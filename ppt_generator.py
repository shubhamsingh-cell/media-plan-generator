#!/usr/bin/env python3
"""
Joveo-branded PowerPoint generator for AI Media Planner.

Generates a polished, data-driven .pptx presentation using python-pptx.
Uses Joveo brand identity: Port Gore navy (#202058), Blue Violet purple (#5A54BE),
Downy teal (#6BB5CE), Light Purple (#8680D6), Light Teal (#A8D8EA), Magenta (#B7669E).
Fonts: Poppins (headings), Inter (body). Incorporates hero stats, section
dividers, quality outcomes grids, channel attribution diagrams, and comparison panels.

Note: This module does not directly import data_orchestrator.py. It receives
orchestrated/enriched data transitively via app.py, which calls the orchestrator
and passes the enriched results into the PPT generation functions.
"""

import io
import json
import logging
import math
import os
import re
import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

logger = logging.getLogger(__name__)

from shared_utils import (
    parse_budget_display,
    INDUSTRY_LABEL_MAP as _SHARED_INDUSTRY_LABEL_MAP,
    internal_qc_mode as _internal_qc_mode,
)

from joveo_brand_2026 import (
    LAVENDER_50 as _LAVENDER_50_HEX,
    LAVENDER_100 as _LAVENDER_100_HEX,
    BLUE_50 as _BLUE_50_HEX,
)

try:
    import plan_currency as _plan_currency
except ImportError:  # pragma: no cover - plan_currency ships with the repo
    _plan_currency = None

import plan_geo as _plan_geo
import display_format as _fmt
import insight_composer as _insight
import gold_standard as _gold_standard

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

try:
    import research
except ImportError:
    research = None

try:
    import trend_engine as _trend_engine

    _HAS_TREND_ENGINE = True
except ImportError:
    _HAS_TREND_ENGINE = False

try:
    import collar_intelligence as _collar_intel

    _HAS_COLLAR_INTEL = True
except ImportError:
    _HAS_COLLAR_INTEL = False

try:
    import matplotlib

    matplotlib.use("Agg")  # Non-interactive backend for server use
    import matplotlib.pyplot as plt
    import matplotlib.patches as mpatches

    _HAS_MATPLOTLIB = True
except ImportError:
    _HAS_MATPLOTLIB = False


# ---------------------------------------------------------------------------
# Chart font (S89): matplotlib defaults to DejaVu Sans, which clashes with the
# Joveo deck (Poppins headings). Bundle Poppins .ttf files (OFL-licensed) under
# fonts/ and register them with matplotlib.font_manager at import so chart PNGs
# render in Poppins. Falls back to DejaVu gracefully when the fonts are absent.
# ---------------------------------------------------------------------------
_CHART_FONT_FAMILY = "DejaVu Sans"  # safe fallback that ships with matplotlib
_FONTS_DIR = Path(__file__).resolve().parent / "fonts"


def _register_chart_fonts() -> str:
    """Register bundled Poppins fonts with matplotlib; return the family to use.

    Best-effort: any missing file or registration error leaves the chart font on
    the DejaVu fallback so chart generation never breaks.
    """
    if not _HAS_MATPLOTLIB:
        return "DejaVu Sans"
    try:
        from matplotlib import font_manager as _fm

        registered = False
        for _ttf in (
            "Poppins-Regular.ttf",
            "Poppins-SemiBold.ttf",
            "Poppins-Bold.ttf",
        ):
            _path = _FONTS_DIR / _ttf
            if _path.is_file():
                try:
                    _fm.fontManager.addfont(str(_path))
                    registered = True
                except Exception as _font_exc:  # noqa: BLE001
                    logger.debug("Poppins font register skipped (%s): %s", _ttf, _font_exc)
        if not registered:
            return "DejaVu Sans"
        # Confirm matplotlib can actually resolve the family before adopting it.
        try:
            resolved = _fm.findfont(
                _fm.FontProperties(family="Poppins"), fallback_to_default=False
            )
            if resolved and "poppins" in resolved.lower():
                plt.rcParams["font.family"] = "Poppins"
                # Keep DejaVu in the sans-serif chain for glyph coverage (e.g.
                # currency symbols Poppins may lack).
                plt.rcParams["font.sans-serif"] = [
                    "Poppins",
                    "DejaVu Sans",
                    "Arial",
                    "sans-serif",
                ]
                return "Poppins"
        except Exception as _resolve_exc:  # noqa: BLE001
            logger.debug("Poppins font resolve failed: %s", _resolve_exc)
    except Exception as _fm_exc:  # noqa: BLE001
        logger.debug("Chart font registration skipped: %s", _fm_exc)
    return "DejaVu Sans"


_CHART_FONT_FAMILY = _register_chart_fonts()


# ---------------------------------------------------------------------------
# Font embedding -- so the deck renders in Poppins even on machines that don't
# have it installed. Without this, PowerPoint/Keynote/Quick Look substitute a
# default (often Times serif on macOS), which is the single biggest reason a
# correctly-built deck can still "look wrong" to a viewer.
# ---------------------------------------------------------------------------
# (typeface, embeddedFont slot, filename). PowerPoint matches an embedded font to
# text runs by typeface name; the deck sets every run to "Poppins".
_EMBED_FONTS = [
    ("Poppins", "regular", "Poppins-Regular.ttf"),
    ("Poppins", "bold", "Poppins-Bold.ttf"),
]
_FONT_REL_TYPE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/font"
)


def _embed_fonts_in_pptx(pptx_bytes: bytes) -> bytes:
    """Inject the bundled Poppins .ttf files into a .pptx so it renders on-brand
    everywhere. Pure zip/OOXML surgery (python-pptx has no font-embed API).

    Adds: ppt/fonts/fontN.fntdata parts, <p:embeddedFontLst> in presentation.xml
    (with embedTrueTypeFonts="1"), font relationships, and the fntdata content
    type. Returns the original bytes unchanged on any error — never break the
    deck over a cosmetic enhancement.
    """
    import zipfile

    fonts = []
    for typeface, slot, fname in _EMBED_FONTS:
        fpath = _FONTS_DIR / fname
        if fpath.exists():
            try:
                fonts.append((typeface, slot, fpath.read_bytes()))
            except OSError as exc:
                logger.debug("Font embed: could not read %s: %s", fname, exc)
    if not fonts:
        return pptx_bytes

    try:
        with zipfile.ZipFile(io.BytesIO(pptx_bytes), "r") as zin:
            pres_xml = zin.read("ppt/presentation.xml").decode("utf-8")
            rels_xml = zin.read("ppt/_rels/presentation.xml.rels").decode("utf-8")
            ct_xml = zin.read("[Content_Types].xml").decode("utf-8")
            if "embedTrueTypeFonts" in pres_xml:
                return pptx_bytes  # already embedded — idempotent

            font_parts = []
            rel_entries = []
            by_face: Dict[str, list] = {}
            for i, (typeface, slot, data) in enumerate(fonts, start=1):
                rid = f"rIdFont{i}"
                partname = f"ppt/fonts/font{i}.fntdata"
                font_parts.append((partname, data))
                rel_entries.append(
                    f'<Relationship Id="{rid}" Type="{_FONT_REL_TYPE}" '
                    f'Target="fonts/font{i}.fntdata"/>'
                )
                by_face.setdefault(typeface, []).append((slot, rid))

            embed_lst = "<p:embeddedFontLst>"
            for face, slots in by_face.items():
                embed_lst += f'<p:embeddedFont><p:font typeface="{face}"/>'
                for slot, rid in slots:
                    embed_lst += f'<p:{slot} r:id="{rid}"/>'
                embed_lst += "</p:embeddedFont>"
            embed_lst += "</p:embeddedFontLst>"

            # presentation.xml: flag + insert embeddedFontLst after notesSz
            # (schema order: ... sldSz, notesSz, embeddedFontLst, defaultTextStyle)
            pres_xml = re.sub(
                r"<p:presentation ",
                '<p:presentation embedTrueTypeFonts="1" ',
                pres_xml,
                count=1,
            )
            m = re.search(r"<p:notesSz[^>]*/>", pres_xml)
            if m:
                pres_xml = pres_xml[: m.end()] + embed_lst + pres_xml[m.end():]
            elif "<p:defaultTextStyle" in pres_xml:
                pres_xml = pres_xml.replace(
                    "<p:defaultTextStyle", embed_lst + "<p:defaultTextStyle", 1
                )
            else:
                pres_xml = pres_xml.replace(
                    "</p:presentation>", embed_lst + "</p:presentation>", 1
                )

            rels_xml = rels_xml.replace(
                "</Relationships>", "".join(rel_entries) + "</Relationships>"
            )
            if "fntdata" not in ct_xml:
                ct_xml = ct_xml.replace(
                    "</Types>",
                    '<Default Extension="fntdata" '
                    'ContentType="application/x-fontdata"/></Types>',
                )

            out = io.BytesIO()
            with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as zout:
                for item in zin.infolist():
                    data = zin.read(item.filename)
                    if item.filename == "ppt/presentation.xml":
                        data = pres_xml.encode("utf-8")
                    elif item.filename == "ppt/_rels/presentation.xml.rels":
                        data = rels_xml.encode("utf-8")
                    elif item.filename == "[Content_Types].xml":
                        data = ct_xml.encode("utf-8")
                    zout.writestr(item, data)
                for partname, data in font_parts:
                    zout.writestr(partname, data)
        return out.getvalue()
    except (KeyError, zipfile.BadZipFile, ValueError) as exc:
        logger.warning("Font embedding skipped (non-fatal): %s", exc)
        return pptx_bytes


def _fix_pptx_package_hygiene(pptx_bytes: bytes, n_slides: int) -> bytes:
    """Post-process package-level metadata python-pptx's default template
    otherwise leaves stale/fabricated (craft:both#3 / craft:both#6):

    - ppt/theme/theme1.xml: the theme's default major/minor Latin typeface
      is the template's stock "Calibri" -- any shape that falls back to the
      theme default (a new text box, PowerPoint's own re-layout) would
      render off-brand. Set it to Poppins, the family this deck actually
      embeds and uses on every slide.
    - docProps/app.xml: the stock template reports 0 slides, 0 paragraphs,
      and "Microsoft Macintosh PowerPoint" v14 as the generating
      application -- nonsense provenance for an 11+-slide AI-generated deck.
      Set the real slide count and the actual generating application.

    Pure zip/OOXML surgery (python-pptx exposes neither). Returns the
    original bytes unchanged on any error -- never break the deck over a
    metadata cleanup.
    """
    import zipfile

    try:
        with zipfile.ZipFile(io.BytesIO(pptx_bytes), "r") as zin:
            names = set(zin.namelist())
            theme_xml = (
                zin.read("ppt/theme/theme1.xml").decode("utf-8")
                if "ppt/theme/theme1.xml" in names
                else None
            )
            app_xml = (
                zin.read("docProps/app.xml").decode("utf-8")
                if "docProps/app.xml" in names
                else None
            )
            if theme_xml is None and app_xml is None:
                return pptx_bytes

            if theme_xml is not None:
                theme_xml = theme_xml.replace(
                    'typeface="Calibri"', 'typeface="Poppins"'
                )

            if app_xml is not None:
                app_xml = re.sub(
                    r"<Application>[^<]*</Application>",
                    "<Application>Nova AI Suite</Application>",
                    app_xml,
                )
                app_xml = re.sub(
                    r"<AppVersion>[^<]*</AppVersion>",
                    "<AppVersion>2026.0</AppVersion>",
                    app_xml,
                )
                app_xml = re.sub(
                    r"<Slides>\d*</Slides>",
                    f"<Slides>{int(n_slides)}</Slides>",
                    app_xml,
                )
                app_xml = re.sub(
                    r"<PresentationFormat>[^<]*</PresentationFormat>",
                    "<PresentationFormat>On-screen Show (16:9)</PresentationFormat>",
                    app_xml,
                )

            out = io.BytesIO()
            with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as zout:
                for item in zin.infolist():
                    data = zin.read(item.filename)
                    if item.filename == "ppt/theme/theme1.xml" and theme_xml is not None:
                        data = theme_xml.encode("utf-8")
                    elif item.filename == "docProps/app.xml" and app_xml is not None:
                        data = app_xml.encode("utf-8")
                    zout.writestr(item, data)
        return out.getvalue()
    except (KeyError, zipfile.BadZipFile, ValueError) as exc:
        logger.warning("Package hygiene fix skipped (non-fatal): %s", exc)
        return pptx_bytes


# ---------------------------------------------------------------------------
# Chart color palette (hex strings for matplotlib, matching Joveo brand)
# ---------------------------------------------------------------------------
# S89: pure Joveo deck dataviz sequence (no off-palette greens/teals). Order
# follows the brand guideline series: purple -> teal -> magenta -> purple-light,
# then deepen within the family. Categorical, high-contrast between neighbors.
_CHART_COLORS = [
    "#5A54BE",  # PURPLE — series 1 (primary accent)
    "#6BB5CE",  # TEAL — series 2
    "#B7669E",  # MAGENTA — series 3
    "#8680D6",  # PURPLE_LIGHT — series 4
    "#202058",  # INDIGO — series 5 (deep)
    "#3E8FAB",  # TEAL deep — series 6
    "#C98BB6",  # MAGENTA light — series 7
    "#3F3A8E",  # PURPLE deep — series 8
]


def _generate_pie_chart_image(labels: List[str], sizes: List[float]) -> Optional[bytes]:
    """Generate a budget allocation pie chart as PNG bytes using matplotlib.

    Args:
        labels: Channel names.
        sizes: Percentage allocations (should sum to ~100).

    Returns:
        PNG image bytes, or None if matplotlib is unavailable or chart fails.
    """
    if not _HAS_MATPLOTLIB or not labels or not sizes:
        return None
    try:
        # S89: Cap slice count so the chart stays legible. Plans with many
        # channels (10+) produced sliver slices and an overflowing legend;
        # keep the top 7 by share and roll the rest into a single "Other".
        _MAX_SLICES = 8
        if len(labels) > _MAX_SLICES:
            pairs = sorted(zip(labels, sizes), key=lambda p: p[1], reverse=True)
            top = pairs[: _MAX_SLICES - 1]
            rest = pairs[_MAX_SLICES - 1 :]
            other_size = sum(sz for _, sz in rest)
            labels = [lbl for lbl, _ in top] + [f"Other ({len(rest)} channels)"]
            sizes = [sz for _, sz in top] + [round(other_size, 1)]

        fig, ax = plt.subplots(figsize=(7, 5), dpi=150)
        fig.patch.set_facecolor("#FFFCF9")

        colors = _CHART_COLORS[: len(labels)]
        # Extend colors if we have more channels than palette entries
        while len(colors) < len(labels):
            colors.append(_CHART_COLORS[len(colors) % len(_CHART_COLORS)])

        wedges, texts, autotexts = ax.pie(
            sizes,
            labels=None,
            autopct=lambda pct: f"{pct:.1f}%" if pct >= 3 else "",
            startangle=90,
            colors=colors,
            pctdistance=0.75,
            wedgeprops={"linewidth": 1.5, "edgecolor": "white"},
        )

        for autotext in autotexts:
            autotext.set_fontsize(9)
            autotext.set_fontweight("bold")
            autotext.set_color("white")

        # Add legend to the right
        legend_labels = [f"{lbl} ({sz:.0f}%)" for lbl, sz in zip(labels, sizes)]
        ax.legend(
            wedges,
            legend_labels,
            title="Channels",
            loc="center left",
            bbox_to_anchor=(1.0, 0.5),
            fontsize=8,
            title_fontsize=9,
            frameon=False,
        )

        ax.set_title(
            "Budget Allocation by Channel",
            fontsize=13,
            fontweight="bold",
            color="#202058",
            pad=15,
        )

        plt.tight_layout()
        buf = io.BytesIO()
        fig.savefig(
            buf, format="png", bbox_inches="tight", facecolor=fig.get_facecolor()
        )
        plt.close(fig)
        buf.seek(0)
        return buf.getvalue()
    except Exception as exc:
        logger.warning("Pie chart generation failed (non-fatal): %s", exc)
        if "fig" in dir():
            try:
                plt.close(fig)
            except Exception:
                pass
        return None


def _generate_funnel_chart_image(
    impressions: int,
    clicks: int,
    applications: int,
    hires: int,
) -> Optional[bytes]:
    """Generate a horizontal funnel chart as PNG bytes using matplotlib.

    Shows the conversion funnel: Impressions -> Clicks -> Applications -> Hires.

    Args:
        impressions: Total impressions.
        clicks: Total clicks.
        applications: Total applications.
        hires: Total hires.

    Returns:
        PNG image bytes, or None if matplotlib is unavailable or chart fails.
    """
    if not _HAS_MATPLOTLIB:
        return None
    try:
        stages = ["Impressions", "Clicks", "Applications", "Hires"]
        values = [
            max(impressions, 1),
            max(clicks, 1),
            max(applications, 1),
            max(hires, 1),
        ]

        fig, ax = plt.subplots(figsize=(8, 4.5), dpi=150)
        fig.patch.set_facecolor("#FFFCF9")

        # S89: brand-only funnel (was a non-deck green #338721 on the last stage)
        funnel_colors = ["#202058", "#5A54BE", "#6BB5CE", "#B7669E"]
        max_val = values[0]

        y_positions = [3.0, 2.0, 1.0, 0.0]
        bar_height = 0.65

        for i, (stage, val, color) in enumerate(zip(stages, values, funnel_colors)):
            width = max(val / max_val, 0.04)  # Minimum visible width
            # Center-align bars
            left = (1.0 - width) / 2.0
            bar = ax.barh(
                y_positions[i],
                width,
                left=left,
                height=bar_height,
                color=color,
                edgecolor="white",
                linewidth=1.5,
                zorder=2,
            )

            # Stage label on the left
            ax.text(
                0.0,
                y_positions[i],
                stage,
                ha="right",
                va="center",
                fontsize=10,
                fontweight="bold",
                color="#202058",
                transform=ax.get_yaxis_transform(),
            )

            # Value label centered in bar
            display_val = f"{val:,}"
            ax.text(
                0.5,
                y_positions[i],
                display_val,
                ha="center",
                va="center",
                fontsize=10,
                fontweight="bold",
                color="white",
                zorder=3,
            )

            # Conversion rate arrow between stages
            if i < len(values) - 1:
                conv_rate = values[i + 1] / values[i] * 100
                ax.annotate(
                    f"{conv_rate:.1f}%",
                    xy=(0.5, y_positions[i] - bar_height / 2 - 0.05),
                    ha="center",
                    va="top",
                    fontsize=8,
                    color="#596780",
                    fontstyle="italic",
                )

        ax.set_xlim(0, 1.0)
        ax.set_ylim(-0.6, 3.8)
        ax.axis("off")

        ax.set_title(
            "Recruitment Conversion Funnel",
            fontsize=13,
            fontweight="bold",
            color="#202058",
            pad=15,
        )

        plt.tight_layout()
        buf = io.BytesIO()
        fig.savefig(
            buf, format="png", bbox_inches="tight", facecolor=fig.get_facecolor()
        )
        plt.close(fig)
        buf.seek(0)
        return buf.getvalue()
    except Exception as exc:
        logger.warning("Funnel chart generation failed (non-fatal): %s", exc)
        if "fig" in dir():
            try:
                plt.close(fig)
            except Exception:
                pass
        return None


# ---------------------------------------------------------------------------
# Brand name casing -- preserves known brand names when title-casing client
# ---------------------------------------------------------------------------
_BRAND_CASING: dict[str, str] = {
    "fedex": "FedEx",
    "linkedin": "LinkedIn",
    "youtube": "YouTube",
    "ibm": "IBM",
    "ups": "UPS",
    "jpmorgan": "JPMorgan",
    "walmart": "Walmart",
    "mcdonalds": "McDonald's",
    "at&t": "AT&T",
    "bmw": "BMW",
    "dhl": "DHL",
    "usps": "USPS",
    "xpo": "XPO",
    "jb hunt": "J.B. Hunt",
    "j.b. hunt": "J.B. Hunt",
    "hca": "HCA",
    "cvs": "CVS",
    "ge": "GE",
    "3m": "3M",
    "bp": "BP",
    "ihg": "IHG",
}


def _proper_client_name(name: str) -> str:
    """Client-facing casing for a client name.

    Known brand overrides (_BRAND_CASING: "AT&T", "J.B. Hunt", ...) win
    first; otherwise defers to display_format.client_display_name's
    word-wise casing. The previous "title-case only if fully upper OR fully
    lower, else leave alone" rule silently passed mixed-case raw input
    straight through -- "atria Senior living" (as literally received) never
    became "Atria Senior Living" because it wasn't ALL lower or ALL upper.
    """
    if not name or name == "Client":
        return name
    lower = name.strip().lower()
    if lower in _BRAND_CASING:
        return _BRAND_CASING[lower]
    return _fmt.client_display_name(name) or name.strip()


# ---------------------------------------------------------------------------
# Constants & Color Palette (Joveo brand identity)
# Primary: Port Gore #202058  |  Accent: Blue Violet #5A54BE
# Secondary: Downy Teal #6BB5CE  |  Extended: Light Purple #8680D6
# Extended: Light Teal #A8D8EA  |  Emphasis: Magenta #B7669E
# ---------------------------------------------------------------------------

# -- Joveo brand primaries --
NAVY = RGBColor(0x20, 0x20, 0x58)  # Port Gore — primary dark / headings
BLUE = RGBColor(0x5A, 0x54, 0xBE)  # Blue Violet — primary accent
MEDIUM_BLUE = RGBColor(0x48, 0x43, 0x9E)  # Deeper purple accent
LIGHT_BLUE = RGBColor(0xDD, 0xDB, 0xFF)  # Light purple background
PALE_BLUE = RGBColor(0xB8, 0xB4, 0xF7)  # Medium purple accent fill
SKY_BLUE = RGBColor(0xA8, 0xD8, 0xEA)  # Light teal (Joveo extended)

# -- Joveo secondary --
TEAL = RGBColor(0x6B, 0xB5, 0xCE)  # Downy Teal — secondary accent
LIGHT_TEAL = RGBColor(0xA8, 0xD8, 0xEA)  # Light teal (Joveo extended)
PALE_TEAL = RGBColor(0xDA, 0xF5, 0xFF)  # Pale teal background

# -- Neutrals --
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xFF, 0xFC, 0xF9)  # Warm white canvas (page bg)
WARM_WHITE = RGBColor(0xFF, 0xFC, 0xF9)  # Card backgrounds
WARM_GRAY = RGBColor(0xEB, 0xE6, 0xE0)  # Borders, dividers
MEDIUM_GRAY = RGBColor(0xD6, 0xCF, 0xC2)  # Subtle separators

# -- Text colors --
DARK_TEXT = RGBColor(0x20, 0x20, 0x58)  # Port Gore for body text
MUTED_TEXT = RGBColor(0x59, 0x67, 0x80)  # Secondary text
LIGHT_MUTED = RGBColor(0x8C, 0x96, 0xA8)  # Tertiary text

# -- Semantic colors --
GREEN = RGBColor(0x33, 0x87, 0x21)  # Positive / beating benchmark
LIGHT_GREEN = RGBColor(0xE6, 0xF2, 0xE0)  # Green background
AMBER = RGBColor(0x3E, 0x8F, 0xAB)  # Teal deep — trailing benchmark
LIGHT_AMBER = RGBColor(0xFD, 0xDB, 0xB2)  # Light bronze background
RED_ACCENT = RGBColor(0xB7, 0x66, 0x9E)  # Magenta — underperformance accent
GOLD = RGBColor(0xB7, 0x66, 0x9E)  # Magenta — emphasis / highlights

# -- Joveo extended palette --
JOVEO_LIGHT_PURPLE = RGBColor(0x86, 0x80, 0xD6)  # Light purple accent
JOVEO_BRONZE = RGBColor(0x3E, 0x8F, 0xAB)  # Teal deep accent / CTAs
JOVEO_PINK = RGBColor(0xB7, 0x66, 0x9E)  # Magenta accent

# -- Deck narrative surfaces (canonical hexes imported from joveo_brand_2026) --
LAVENDER_50 = RGBColor.from_string(_LAVENDER_50_HEX.lstrip("#"))  # zebra rows / strips
LAVENDER_100 = RGBColor.from_string(_LAVENDER_100_HEX.lstrip("#"))  # Push card surface
BLUE_50 = RGBColor.from_string(_BLUE_50_HEX.lstrip("#"))  # Pull card surface

# -- Fonts (Poppins headings + body -- Joveo deck 2026) --
# craft:both#5: "Inter" body-font runs were never embedded in the PPTX (only
# Poppins ships in ppt/fonts/ + presentation.xml's embeddedFontLst), so any
# machine without Inter installed silently substituted a fallback font.
# Standardize every text run onto the one family this deck actually embeds
# (the fix instructions' own listed alternative to shipping an Inter .ttf,
# which isn't present in fonts/ to embed).
FONT_FAMILY = "Poppins"  # Brand heading / title font
FONT_BODY = "Poppins"  # Brand body font (was "Inter" -- not embedded)

# Slide dimensions (16:9 widescreen)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


_AI_TRAINING_PHRASES = (
    "ai training",
    "ai trainer",
    "data labeling",
    "data annotation",
    "data label",
    "rlhf",
    "llm training",
)
_AI_TRAINING_TOKENS = frozenset({"ai", "annotator", "annotation"})
_AI_TRAINING_TOKEN_RE = re.compile(r"[a-z0-9]+")


def _is_ai_training_plan(plan_data: Dict[str, Any]) -> bool:
    """Return True if the plan's industry/roles match an AI-training engagement.

    CPA-reference slides carry AI-trainer-specific role data and must not be
    shown to unrelated clients (healthcare, logistics, etc.). Matches on
    token/phrase boundaries only -- a raw substring scan previously matched
    "ai" inside "Maintenance Technician" and "Supply Chain" false-positived
    the same way, misfiring the AI-training slides for unrelated briefs.
    """
    if not isinstance(plan_data, dict):
        return False
    parts: List[str] = []
    industry = plan_data.get("industry")
    if industry:
        parts.append(str(industry))
    roles = plan_data.get("roles") or []
    if isinstance(roles, list):
        for r in roles:
            if isinstance(r, dict):
                parts.append(str(r.get("name") or r.get("title") or ""))
            else:
                parts.append(str(r))
    elif roles:
        parts.append(str(roles))
    text = " ".join(parts).lower()
    if any(phrase in text for phrase in _AI_TRAINING_PHRASES):
        return True
    tokens = set(_AI_TRAINING_TOKEN_RE.findall(text))
    return bool(tokens & _AI_TRAINING_TOKENS)


def _kb_content_key(data: Dict[str, Any]) -> str:
    """Resolve the industry key that selects this plan's KB-sourced deck
    content (data/joveo_media_plan_deck_2026.json v2.0's ``cpa_reference`` /
    ``case_study``, both industry-keyed with only ``'ai_training'``
    populated -- see the KB's ``_meta`` block). AI-training plans always
    resolve to ``'ai_training'`` regardless of the plan's raw industry
    string; every other plan resolves to its own industry key, which will
    not be present in either KB section (by design -- fabricating CPA
    ranges or case-study proof for unrelated industries is out of bounds),
    so callers must treat a missing key as "drop this content"."""
    plan_gate = {"industry": data.get("industry"), "roles": data.get("roles")}
    if _is_ai_training_plan(plan_gate):
        return "ai_training"
    return str(data.get("industry") or "").strip()


def _load_deck_kb() -> Dict[str, Any]:
    """Load the Joveo media-plan deck content (methodology, push/pull, CPA
    reference, case study, next steps) from the KB.

    Mirrors pdf_generator._load_deck_kb. Returns ``{}`` if the file is
    unavailable so the presentation degrades gracefully.
    """
    try:
        path = Path(__file__).parent / "data" / "joveo_media_plan_deck_2026.json"
        with open(path, "r", encoding="utf-8") as fh:
            return json.load(fh)
    except (OSError, ValueError) as exc:  # missing file / malformed JSON
        logger.warning("deck KB unavailable, skipping narrative slides: %s", exc)
        return {}


def _set_body_font(tf) -> None:
    """Switch every run in a text frame to the Inter body font."""
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.name = FONT_BODY


# fix/plan-quality-8 (prod Atria bundle): source text (company bios pulled
# from enrichment APIs) sometimes already carries its own trailing "..." /
# "…" from an upstream snippet truncation. Re-truncating that raw text
# without stripping the existing marker first produced a doubled ellipsis,
# or -- once the marker fell inside our own maxlen window -- a cut that
# read as if it landed mid-word. Strip any pre-existing trailing ellipsis
# before deciding whether/where to cut so the final text always ends in
# exactly one clean ellipsis, never a partial word.
_TRAILING_ELLIPSIS_RE = re.compile(r"(?:\.{2,}|…)\s*$")


def _trunc_word(s: str, maxlen: int = 500) -> str:
    """Truncate text at a word boundary to prevent PPT text box overflow.

    Always produces a single clean trailing ellipsis and never cuts inside
    a word -- even when the input already carries its own raw/partial
    ellipsis (e.g. a pre-truncated upstream API description).
    """
    s = _TRAILING_ELLIPSIS_RE.sub("", str(s or "")).rstrip()
    if len(s) <= maxlen:
        return s
    cut = s[:maxlen].rsplit(" ", 1)[0].rstrip(" ,;:-")
    if not cut:
        # No space anywhere in the window (one giant unbroken token) --
        # fall back to a hard cut rather than returning nothing.
        cut = s[:maxlen].rstrip(" ,;:-")
    return cut + "..."


# ---------------------------------------------------------------------------
# Industry Benchmark Data
# NOTE: Canonical benchmark source is trend_engine.py. These values are fallbacks only.
# See trend_engine.get_benchmark() for authoritative CPC/CPA/CPM data with
# seasonal, regional, and collar-type adjustments. The _get_benchmarks()
# function below attempts to use trend_engine first, falling back to these.
# ---------------------------------------------------------------------------

BENCHMARKS: Dict[str, Dict[str, str]] = {
    "healthcare_medical": {
        "cpa": "$35 - $85",
        "cpc": "$0.90 - $3.50",
        "cph": "$9K - $12K",
        "apply_rate": "3.2% - 4.5%",
    },
    "tech_engineering": {
        "cpa": "$25 - $75",
        "cpc": "$1.20 - $4.50",
        "cph": "$6K - $22K",
        "apply_rate": "6.41%",
    },
    "retail_consumer": {
        "cpa": "$8 - $21",
        "cpc": "$0.25 - $1.00",
        "cph": "$2.7K - $4K",
        "apply_rate": "4.5% - 5.8%",
    },
    "general_entry_level": {
        "cpa": "$10 - $25",
        "cpc": "$0.35 - $1.30",
        "cph": "$2K - $4.7K",
        "apply_rate": "5.5% - 6.1%",
    },
    "finance_banking": {
        "cpa": "$21 - $65",
        "cpc": "$0.90 - $3.50",
        "cph": "$5K - $12K",
        "apply_rate": "5.0% - 6.0%",
    },
    "logistics_supply_chain": {
        "cpa": "$15 - $52",
        "cpc": "$0.40 - $1.80",
        "cph": "$4.5K - $8K",
        "apply_rate": "4.0% - 5.2%",
    },
    "hospitality_travel": {
        "cpa": "$8 - $25",
        "cpc": "$0.22 - $1.00",
        "cph": "$2.5K - $4K",
        "apply_rate": "4.0% - 5.0%",
    },
    "blue_collar_trades": {
        "cpa": "$12 - $35",
        "cpc": "$0.40 - $1.60",
        "cph": "$3.5K - $5.6K",
        "apply_rate": "4.0% - 5.5%",
    },
    "pharma_biotech": {
        "cpa": "$40 - $110",
        "cpc": "$1.50 - $5.00",
        "cph": "$8K - $18K",
        "apply_rate": "3.8% - 5.2%",
    },
}

# ---------------------------------------------------------------------------
# Industry-Specific Hiring Challenges (Complication column)
# ---------------------------------------------------------------------------

COMPLICATIONS: Dict[str, List[str]] = {
    "healthcare_medical": [
        "Clinical talent shortages persist nationally",
        "CPA exceeds $35+ for standing-up roles",
        # consistency:atria#3: the prior "18% higher churn vs. 2023" claim
        # traced to no field anywhere in the workbook/KB. Replaced with the
        # real, sourced NSI figure the KB actually carries (staff RN
        # turnover 17.6% in 2025, up from 16.4% in 2024) -- see
        # data/healthcare_specialty_pay_2026.json's nursing_demand_2026
        # block, the same source excel_v2's Executive Summary quotes.
        "Staff RN turnover climbed to 17.6% in 2025, up from 16.4% in 2024 (NSI)",
        "Credentialing requirements slow time-to-fill",
    ],
    "tech_engineering": [
        "White-collar recession creating surplus but CPCs remain high",
        "AI/ML roles still command premium sourcing costs",
        "Senior/specialized roles still average 45+ days to fill",
        "Remote-first expectations complicate geo-targeting",
    ],
    "retail_consumer": [
        "64,000 retail jobs shed in 2025 Q1",
        "Retail CPA rising sharply despite market softening",
        "5,800+ store closures accelerating talent displacement",
        "Seasonal demand creates volatile cost spikes",
    ],
    "hospitality_travel": [
        "Hospitality CPA surging as travel demand outpaces labor supply",
        "Extreme seasonal demand swings in Q2/Q4",
        "High turnover-driven churn exceeds 73%",
        "Hourly wage competition from adjacent industries",
    ],
    "logistics_supply_chain": [
        "Logistics CPA rising significantly as warehouse/CDL demand grows",
        # strategy:manpower#2: "$52+ CPA" traced to no field anywhere in the
        # workbook/KB and was invented precision. The real, sourced
        # cdl_drivers CPA benchmark (same field excel_v2's Recruitment
        # Benchmarks table quotes) is $25-$50 -- see
        # data/recruitment_benchmarks_deep.json
        # industry_benchmarks.logistics_supply_chain.cpa.cdl_drivers.
        "CDL/last-mile roles most expensive within the range, at $25-$50 CPA",
        "Automation creating new hybrid role types",
        "Warehouse labor competing with gig economy",
    ],
    "finance_banking": [
        "Finance CPA climbing as compliance and fintech hiring intensify",
        "Compliance-heavy hiring extends cycles by 2-3 weeks",
        "Extensive background checks inflate cost-per-hire",
        "Fintech competition drawing mid-career talent",
    ],
    "general_entry_level": [
        "CPCs trending upward across entry-level roles",
        "Seasonal Q4 spikes compress planning windows",
        "Apply rates improving but quality remains a challenge",
        "High-volume funnels require aggressive top-of-funnel spend",
    ],
    "blue_collar_trades": [
        "Skilled trades gap widening as workforce ages",
        "CPA elevated for certified/licensed positions",
        "Geographic mismatch between supply and demand",
        "Apprenticeship pipelines insufficient for near-term needs",
    ],
    "pharma_biotech": [
        "Regulatory talent scarcity drives $110+ CPAs",
        "Clinical trial staffing requires hyper-niche sourcing",
        "PhD-level roles average 60+ days to fill",
        "Compliance training costs add $3K-$5K per hire",
    ],
}

# ---------------------------------------------------------------------------
# Default Channel Allocations
# ---------------------------------------------------------------------------

# Default allocation (used as fallback)
CHANNEL_ALLOC: Dict[str, Dict[str, Any]] = {
    "programmatic_dsp": {
        "label": "Programmatic DSP",
        "pct": 35,
        "color": NAVY,
        "category": "Programmatic",
    },
    "global_boards": {
        "label": "Global Job Boards",
        "pct": 20,
        "color": BLUE,
        "category": "Job Boards",
    },
    "niche_boards": {
        "label": "Niche / Industry Boards",
        "pct": 15,
        "color": MEDIUM_BLUE,
        "category": "Job Boards",
    },
    "social_media": {
        "label": "Social Media",
        "pct": 12,
        "color": SKY_BLUE,
        "category": "Social",
    },
    "regional_boards": {
        "label": "Regional Boards",
        "pct": 8,
        "color": PALE_BLUE,
        "category": "Job Boards",
    },
    "employer_branding": {
        "label": "Employer Branding",
        "pct": 5,
        "color": TEAL,
        "category": "Employer Brand",
    },
    "apac_regional": {
        "label": "APAC Regional",
        "pct": 3,
        "color": LIGHT_TEAL,
        "category": "Job Boards",
    },
    "emea_regional": {
        "label": "EMEA Regional",
        "pct": 2,
        "color": PALE_TEAL,
        "category": "Job Boards",
    },
}

# ═══════════════════════════════════════════════════════════════════════════════
# INDUSTRY ALLOCATION PROFILES
# Methodology: Derived from analysis of 200+ recruitment media campaigns across
# 17 industries (2024-2025). Validated against Appcast 2026 Recruitment Marketing
# Benchmark Report, LinkedIn Talent Solutions data, and Joveo first-party campaign
# data (66M+ views, 11M+ clicks across 2 accounts).
# Re-calibration: Quarterly review against live campaign performance data.
# Last validated: Q1 2026
# ═══════════════════════════════════════════════════════════════════════════════

# ── Industry-specific allocation profiles ──
# Each profile shifts percentages to match industry hiring patterns.
# The channel keys match CHANNEL_ALLOC keys; only "pct" differs.
INDUSTRY_ALLOC_PROFILES: Dict[str, Dict[str, int]] = {
    # Healthcare: heavier on niche medical boards, less programmatic
    "healthcare_medical": {
        "programmatic_dsp": 22,
        "global_boards": 15,
        "niche_boards": 30,
        "social_media": 10,
        "regional_boards": 10,
        "employer_branding": 8,
        "apac_regional": 3,
        "emea_regional": 2,
    },
    # Tech: heavier on programmatic/digital and social, moderate niche
    "tech_engineering": {
        "programmatic_dsp": 30,
        "global_boards": 15,
        "niche_boards": 20,
        "social_media": 18,
        "regional_boards": 5,
        "employer_branding": 7,
        "apac_regional": 3,
        "emea_regional": 2,
    },
    # Finance: balanced with strong niche and employer branding
    "finance_banking": {
        "programmatic_dsp": 25,
        "global_boards": 18,
        "niche_boards": 25,
        "social_media": 10,
        "regional_boards": 7,
        "employer_branding": 10,
        "apac_regional": 3,
        "emea_regional": 2,
    },
    # Retail/consumer & hospitality: high-volume programmatic + social
    "retail_consumer": {
        "programmatic_dsp": 38,
        "global_boards": 22,
        "niche_boards": 8,
        "social_media": 20,
        "regional_boards": 7,
        "employer_branding": 3,
        "apac_regional": 1,
        "emea_regional": 1,
    },
    "hospitality_travel": {
        "programmatic_dsp": 38,
        "global_boards": 22,
        "niche_boards": 8,
        "social_media": 20,
        "regional_boards": 7,
        "employer_branding": 3,
        "apac_regional": 1,
        "emea_regional": 1,
    },
    # General / entry-level: programmatic-heavy, broad reach
    "general_entry_level": {
        "programmatic_dsp": 40,
        "global_boards": 22,
        "niche_boards": 8,
        "social_media": 15,
        "regional_boards": 10,
        "employer_branding": 3,
        "apac_regional": 1,
        "emea_regional": 1,
    },
    # Blue-collar/trades: programmatic + regional, less niche digital
    "blue_collar_trades": {
        "programmatic_dsp": 35,
        "global_boards": 20,
        "niche_boards": 10,
        "social_media": 15,
        "regional_boards": 15,
        "employer_branding": 3,
        "apac_regional": 1,
        "emea_regional": 1,
    },
    # Aerospace/defense: niche-heavy, security-cleared boards matter
    "aerospace_defense": {
        "programmatic_dsp": 20,
        "global_boards": 15,
        "niche_boards": 30,
        "social_media": 8,
        "regional_boards": 10,
        "employer_branding": 12,
        "apac_regional": 3,
        "emea_regional": 2,
    },
    # Pharma/biotech: niche scientific boards + employer branding
    "pharma_biotech": {
        "programmatic_dsp": 22,
        "global_boards": 15,
        "niche_boards": 28,
        "social_media": 10,
        "regional_boards": 8,
        "employer_branding": 12,
        "apac_regional": 3,
        "emea_regional": 2,
    },
    # Education: niche-heavy (HigherEdJobs etc.), moderate social
    "education": {
        "programmatic_dsp": 20,
        "global_boards": 18,
        "niche_boards": 28,
        "social_media": 12,
        "regional_boards": 10,
        "employer_branding": 7,
        "apac_regional": 3,
        "emea_regional": 2,
    },
    # Legal services: niche + employer brand focused
    "legal_services": {
        "programmatic_dsp": 22,
        "global_boards": 18,
        "niche_boards": 28,
        "social_media": 8,
        "regional_boards": 8,
        "employer_branding": 11,
        "apac_regional": 3,
        "emea_regional": 2,
    },
    # Manufacturing/automotive: programmatic + regional + niche trade
    "automotive": {
        "programmatic_dsp": 30,
        "global_boards": 18,
        "niche_boards": 18,
        "social_media": 10,
        "regional_boards": 15,
        "employer_branding": 5,
        "apac_regional": 2,
        "emea_regional": 2,
    },
    # Energy/utilities: niche trade boards + regional
    "energy_utilities": {
        "programmatic_dsp": 25,
        "global_boards": 15,
        "niche_boards": 25,
        "social_media": 8,
        "regional_boards": 15,
        "employer_branding": 7,
        "apac_regional": 3,
        "emea_regional": 2,
    },
    # Mental health: niche clinical + employer brand
    "mental_health": {
        "programmatic_dsp": 22,
        "global_boards": 18,
        "niche_boards": 28,
        "social_media": 10,
        "regional_boards": 8,
        "employer_branding": 9,
        "apac_regional": 3,
        "emea_regional": 2,
    },
    # Logistics/supply chain/trucking: niche CDL boards dominate for
    # trucking/transportation (CDLjobs, TruckersReport, DriveMyWay),
    # programmatic for volume, regional for local routes.
    # S49 Issue 15: Raised niche_boards from 12% to 35% for trucking alignment.
    "logistics_supply_chain": {
        "programmatic_dsp": 25,
        "global_boards": 15,
        "niche_boards": 35,  # CDLjobs, TruckersReport, DriveMyWay
        "social_media": 10,
        "regional_boards": 10,
        "employer_branding": 5,
        "apac_regional": 0,
        "emea_regional": 0,
    },
    # Insurance: niche + professional boards
    "insurance": {
        "programmatic_dsp": 25,
        "global_boards": 18,
        "niche_boards": 25,
        "social_media": 10,
        "regional_boards": 7,
        "employer_branding": 10,
        "apac_regional": 3,
        "emea_regional": 2,
    },
    # Maritime/marine: niche-heavy, regional
    "maritime_marine": {
        "programmatic_dsp": 20,
        "global_boards": 15,
        "niche_boards": 30,
        "social_media": 8,
        "regional_boards": 15,
        "employer_branding": 7,
        "apac_regional": 3,
        "emea_regional": 2,
    },
}


def _get_industry_alloc(
    industry: str,
    budget_str: str = "",
    num_roles: int = 0,
    roles: list = None,
    locations: list = None,
) -> Dict[str, Dict[str, Any]]:
    """Return a copy of CHANNEL_ALLOC with percentages adjusted for industry, budget, roles, locations."""
    # Use manual dict copy to avoid deepcopy issues with RGBColor objects
    base = {k: dict(v) for k, v in CHANNEL_ALLOC.items()}

    # Step 1: Apply industry profile
    profile = INDUSTRY_ALLOC_PROFILES.get(industry)
    if profile:
        for key in base:
            if key in profile:
                base[key]["pct"] = profile[key]

    # Step 2: Adjust for budget size
    budget_val = _parse_budget_number(budget_str) if budget_str else None
    if budget_val is not None:
        if budget_val < 50000:
            # Small budget: concentrate on top 3-4 channels, cut low-impact ones
            base["employer_branding"]["pct"] = max(
                1, base["employer_branding"]["pct"] - 3
            )
            base["apac_regional"]["pct"] = max(0, base["apac_regional"]["pct"] - 2)
            base["emea_regional"]["pct"] = max(0, base["emea_regional"]["pct"] - 1)
            base["programmatic_dsp"]["pct"] += 4
            base["global_boards"]["pct"] += 2
        elif budget_val > 500000:
            # Large budget: spread wider, invest in branding
            base["employer_branding"]["pct"] += 4
            base["regional_boards"]["pct"] += 2
            base["social_media"]["pct"] += 2
            base["programmatic_dsp"]["pct"] -= 5
            base["global_boards"]["pct"] -= 3

    # Step 3: Adjust for number of roles (more roles = more diverse mix)
    if num_roles and num_roles > 10:
        base["niche_boards"]["pct"] += 3
        base["regional_boards"]["pct"] += 2
        base["programmatic_dsp"]["pct"] -= 3
        base["global_boards"]["pct"] -= 2

    # Step 4: Adjust for seniority mix (if roles provided)
    if roles:
        _role_strs = [
            (
                r
                if isinstance(r, str)
                else (
                    r.get("title", r.get("role", str(r)))
                    if isinstance(r, dict)
                    else str(r)
                )
            )
            for r in roles
        ]
        roles_lower = " ".join(r.lower() for r in _role_strs)
        senior_keywords = [
            "executive",
            "director",
            "vp",
            "chief",
            "president",
            "c-suite",
            "senior",
            "head of",
            "principal",
            "fellow",
        ]
        junior_keywords = [
            "intern",
            "entry",
            "junior",
            "associate",
            "trainee",
            "assistant",
            "coordinator",
            "clerk",
        ]
        senior_count = sum(1 for kw in senior_keywords if kw in roles_lower)
        junior_count = sum(1 for kw in junior_keywords if kw in roles_lower)

        if senior_count > junior_count:
            # Senior-heavy: more niche/executive boards, more employer branding
            base["niche_boards"]["pct"] += 4
            base["employer_branding"]["pct"] += 3
            base["social_media"]["pct"] -= 3
            base["programmatic_dsp"]["pct"] -= 4
        elif junior_count > senior_count:
            # Junior-heavy: more social, more global boards
            base["social_media"]["pct"] += 5
            base["global_boards"]["pct"] += 3
            base["niche_boards"]["pct"] -= 4
            base["employer_branding"]["pct"] -= 2
            base["programmatic_dsp"]["pct"] -= 2

    # Step 5: Geo-filter -- remove APAC/EMEA for US-only plans
    if locations:
        _locs_list = locations if isinstance(locations, list) else [locations]
        _us_terms = {"us", "usa", "united states", ""}
        _all_us = all(
            str(loc.get("country") if isinstance(loc, dict) else loc).lower().strip()
            in _us_terms
            for loc in _locs_list
        )
        if _all_us:
            for _intl_key in ("apac_regional", "emea_regional"):
                if _intl_key in base:
                    _freed = base[_intl_key]["pct"]
                    base[_intl_key]["pct"] = 0
                    if _freed > 0:
                        _top = max(
                            (k for k in base if k != _intl_key),
                            key=lambda k: base[k]["pct"],
                        )
                        base[_top]["pct"] += _freed

    # Ensure no negative percentages
    for key in base:
        base[key]["pct"] = max(
            0 if key in ("apac_regional", "emea_regional") else 1, base[key]["pct"]
        )

    # Filter out zero-allocation channels
    base = {k: v for k, v in base.items() if v["pct"] > 0}

    # Normalize to 100%
    total = sum(v["pct"] for v in base.values())
    if total > 0 and total != 100:
        for key in base:
            base[key]["pct"] = round(base[key]["pct"] / total * 100)
        diff = 100 - sum(v["pct"] for v in base.values())
        if diff != 0:
            # Add remainder to largest category
            largest = max(base, key=lambda k: base[k]["pct"])
            base[largest]["pct"] += diff

    return base


# Human-readable goal labels
GOAL_LABELS: Dict[str, str] = {
    "brand_awareness": "Brand Awareness",
    "high_volume": "High-Volume Hiring",
    "diversity_hiring": "Diversity & Inclusion",
    "cost_efficiency": "Cost Efficiency",
    "quality_candidates": "Quality Candidates",
    "passive_talent": "Passive Talent Reach",
    "employer_branding": "Employer Branding",
    "retention": "Retention Focus",
    "speed_to_hire": "Speed to Hire",
    "geographic_expansion": "Geographic Expansion",
}

WORK_ENV_LABELS: Dict[str, str] = {
    "hybrid": "Hybrid",
    "remote": "Remote",
    "on_site": "On-Site",
    "on-site": "On-Site",
    "flexible": "Flexible",
}

# Industry benchmark comparison data for side-by-side panel
INDUSTRY_BENCHMARKS_COMPARISON: Dict[str, Dict[str, Any]] = {
    "healthcare_medical": {
        "avg_channels": 4,
        "avg_budget_pct_programmatic": 28,
        "avg_apply_rate": 3.8,
        "avg_time_to_fill": 42,
        "avg_cpa": 60,
        "estimated_reach_multiplier": 1.0,
    },
    "tech_engineering": {
        "avg_channels": 5,
        "avg_budget_pct_programmatic": 32,
        "avg_apply_rate": 6.4,
        "avg_time_to_fill": 35,
        "avg_cpa": 50,
        "estimated_reach_multiplier": 1.1,
    },
    "retail_consumer": {
        "avg_channels": 3,
        "avg_budget_pct_programmatic": 25,
        "avg_apply_rate": 5.1,
        "avg_time_to_fill": 28,
        "avg_cpa": 14,
        "estimated_reach_multiplier": 0.9,
    },
    "general_entry_level": {
        "avg_channels": 4,
        "avg_budget_pct_programmatic": 30,
        "avg_apply_rate": 5.8,
        "avg_time_to_fill": 30,
        "avg_cpa": 18,
        "estimated_reach_multiplier": 1.0,
    },
    "finance_banking": {
        "avg_channels": 4,
        "avg_budget_pct_programmatic": 30,
        "avg_apply_rate": 5.5,
        "avg_time_to_fill": 38,
        "avg_cpa": 43,
        "estimated_reach_multiplier": 1.0,
    },
    "logistics_supply_chain": {
        "avg_channels": 4,
        "avg_budget_pct_programmatic": 30,
        "avg_apply_rate": 4.6,
        "avg_time_to_fill": 32,
        "avg_cpa": 34,
        "estimated_reach_multiplier": 1.0,
    },
    "hospitality_travel": {
        "avg_channels": 3,
        "avg_budget_pct_programmatic": 22,
        "avg_apply_rate": 4.5,
        "avg_time_to_fill": 25,
        "avg_cpa": 16,
        "estimated_reach_multiplier": 0.9,
    },
    "blue_collar_trades": {
        "avg_channels": 3,
        "avg_budget_pct_programmatic": 26,
        "avg_apply_rate": 4.8,
        "avg_time_to_fill": 30,
        "avg_cpa": 24,
        "estimated_reach_multiplier": 0.9,
    },
    "pharma_biotech": {
        "avg_channels": 5,
        "avg_budget_pct_programmatic": 35,
        "avg_apply_rate": 4.5,
        "avg_time_to_fill": 55,
        "avg_cpa": 75,
        "estimated_reach_multiplier": 1.1,
    },
}


# ===================================================================
# Helper utilities
# ===================================================================


# ---------------------------------------------------------------------------
# Number Formatting Helpers
# ---------------------------------------------------------------------------

_CURRENCY_SYMBOLS = {
    "USD": "$",
    "EUR": "\u20ac",
    "GBP": "\u00a3",
    "INR": "\u20b9",
    "JPY": "\u00a5",
    "CNY": "\u00a5",
    "AUD": "A$",
    "CAD": "C$",
    "SGD": "S$",
    "HKD": "HK$",
    "NZD": "NZ$",
    "CHF": "CHF ",
    "SEK": "kr ",
    "NOK": "kr ",
    "DKK": "kr ",
    "BRL": "R$",
    "ZAR": "R ",
    "MXN": "MX$",
    "KRW": "\u20a9",
    "THB": "\u0e3f",
    "MYR": "RM ",
    "PHP": "\u20b1",
    "IDR": "Rp ",
    "AED": "AED ",
    "SAR": "SAR ",
}


# ---------------------------------------------------------------------------
# Plan currency (S89): money figures must render in the plan's own currency, not
# a hardcoded "$". ``generate_pptx`` resolves the plan's ISO code once up front
# and stashes it module-side so the many existing ``_fmt_currency`` / salary /
# budget call sites localize without each having to thread ``data`` through.
# Defaults to USD; planning-math benchmarks that are intentionally USD-coded
# pass ``currency="USD"`` explicitly (those call sites bypass the active
# currency by passing currency="USD").
#
# S89 concurrency fix: the active currency is THREAD-LOCAL, not a module global.
# generate_pptx runs on concurrent per-job daemon threads under the threading
# HTTP server, so a shared global would let a £ plan and a $ plan racing in
# parallel render each other's symbol. threading.local gives each generation
# thread its own value.
# ---------------------------------------------------------------------------
import threading as _threading  # noqa: E402

_currency_tls = _threading.local()


def _get_active_currency() -> str:
    """Active plan currency for THIS thread (defaults to USD)."""
    return getattr(_currency_tls, "code", "USD") or "USD"


def _plan_currency_code(data: Optional[Dict]) -> str:
    """Resolve the ISO currency code for a plan from its data. Defaults to USD.

    Order of precedence:
      1. An explicit ``currency`` / ``currency_code`` on the plan data.
      2. ``plan_currency.currency_for_country`` applied to the plan's locations
         (trailing "City, ST, Country" token) or an explicit country field.
      3. USD.
    Never raises.
    """
    if not isinstance(data, dict):
        return "USD"
    explicit = data.get("currency_code") or data.get("currency")
    if isinstance(explicit, str) and explicit.strip():
        return explicit.strip().upper()
    if _plan_currency is None:
        return "USD"
    candidates: List[str] = []
    for key in ("country", "primary_location"):
        val = data.get(key)
        if isinstance(val, str) and val.strip():
            candidates.append(val)
    locs = data.get("locations") or []
    if isinstance(locs, (list, tuple)):
        for loc in locs:
            if isinstance(loc, str) and loc.strip():
                candidates.append(loc)
            elif isinstance(loc, dict):
                country = loc.get("country") or loc.get("location") or ""
                if isinstance(country, str) and country.strip():
                    candidates.append(country)
    for cand in candidates:
        try:
            code = _plan_currency.currency_for_country(cand)
        except Exception:  # noqa: BLE001 - resolution is best-effort
            code = None
        if code:
            return code
    return "USD"


def _set_active_currency(data: Optional[Dict]) -> str:
    """Resolve and remember the plan currency for the duration of generation.

    Stored thread-locally so concurrent generations don't clobber each other.
    """
    code = _plan_currency_code(data)
    _currency_tls.code = code
    if isinstance(data, dict):
        data["_plan_currency_code"] = code
    return code


def _cur_symbol(currency: Optional[str] = None) -> str:
    """Return the symbol for ``currency`` (or the active plan currency)."""
    code = (currency or _get_active_currency() or "USD").strip().upper()
    if _plan_currency is not None:
        return _plan_currency.symbol_for_code(code)
    return _CURRENCY_SYMBOLS.get(code, "$")


# ===================================================================
# INVISIBLE SCAFFOLD
# Exact reskin of the Joveo "Invisible Media Planning Approach" deck.
# Geometry tokens were extracted from the source .pptx (10 x 5.625 in)
# and scaled x1.3333 onto the 13.333 x 7.5 in generator canvas.
# Every content slide is built from these helpers so the whole deck
# inherits one clean, light, Invisible-matching design language
# (light lavender canvas, top accent rule + wordmark, big indigo title,
# accent-bar cards, lavender question band, dark takeaway callout).
# ===================================================================
INV_ACCENTS = [BLUE, TEAL, JOVEO_PINK, JOVEO_LIGHT_PURPLE, AMBER]  # purple/teal/magenta/lilac/teal-deep
INK_SLATE = RGBColor(0x33, 0x33, 0x4F)  # card body text (source #33334E)
INV_CANVAS = LAVENDER_50  # #F4F4FF cool light page canvas
INV_PILL = LAVENDER_50    # #F4F4FF soft pill / row surface
INV_QBAND = LAVENDER_100  # #ECEAF7 question-band surface
INV_CONTENT_TOP = 1.55    # inches: y where slide body starts (below title)


def _inv_canvas(slide, wordmark: bool = True):
    """Paint the light Invisible page: lavender canvas, top accent rule, wordmark."""
    _add_filled_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, INV_CANVAS)
    # Full-width top accent rule (purple-light) -- source 0.08in x1.333
    _add_filled_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.10), JOVEO_LIGHT_PURPLE)
    if wordmark:
        _add_textbox(
            slide, Inches(0.65), Inches(0.27), Inches(2.0), Inches(0.4),
            text="joveo", font_size=19, bold=True, color=BLUE,
        )


def _inv_header(slide, title: str, subtitle: str = "", wordmark: bool = True) -> float:
    """Invisible content-slide header (canvas + big indigo title + optional subtitle).

    Returns the y (in inches) where slide body content should start.
    """
    _inv_canvas(slide, wordmark=wordmark)
    _add_textbox(
        slide, Inches(0.63), Inches(0.72), Inches(12.05), Inches(0.64),
        text=title, font_size=30, bold=True, color=NAVY,
    )
    if subtitle:
        _add_textbox(
            slide, Inches(0.65), Inches(1.38), Inches(12.05), Inches(0.36),
            text=subtitle, font_size=14, italic=True, color=BLUE,
        )
        return 1.92
    return INV_CONTENT_TOP


def _inv_question_band(slide, qtext: str, top: float = 1.5) -> float:
    """Lavender question band with a purple 'Q' chip (Q1-Q4 framing slides).

    Returns the y (inches) just below the band.
    """
    y = Inches(top)
    h = Inches(0.86)
    _add_rounded_rect(slide, Inches(0.65), y, Inches(12.1), h, INV_QBAND)
    _add_rounded_rect(slide, Inches(0.65), y, Inches(0.72), h, BLUE)
    _add_textbox(
        slide, Inches(0.65), y, Inches(0.72), h, text="Q", font_size=18, bold=True,
        color=WHITE, alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    _add_textbox(
        slide, Inches(1.6), y, Inches(10.9), h, text=qtext, font_size=13, italic=True,
        color=NAVY, anchor=MSO_ANCHOR.MIDDLE,
    )
    return top + 0.86 + 0.18


def _inv_card(
    slide, x, y, w, h, header: str, body: str = "", accent=None,
    header_color=None, header_size: int = 14, body_size: int = 11,
):
    """White rounded card + colored top accent bar + bold header + body text."""
    accent = accent if accent is not None else BLUE
    card = _add_rounded_rect(slide, x, y, w, h, WHITE)
    try:
        card.line.color.rgb = WARM_GRAY
        card.line.width = Pt(0.75)
    except Exception:
        pass
    _add_filled_rect(slide, x, y, w, Inches(0.08), accent)
    pad = Inches(0.24)
    _add_textbox(
        slide, x + pad, y + Inches(0.2), w - pad * 2, Inches(0.46),
        text=header, font_size=header_size, bold=True, color=header_color or accent,
    )
    if body:
        _add_textbox(
            slide, x + pad, y + Inches(0.64), w - pad * 2, h - Inches(0.78),
            text=body, font_size=body_size, color=INK_SLATE,
        )
    return card


def _inv_callout(slide, text: str, top: float = 6.5):
    """Dark-indigo takeaway callout bar near the bottom (white text), full width."""
    y = Inches(top)
    _add_rounded_rect(slide, Inches(0.65), y, Inches(12.05), Inches(0.52), NAVY)
    _add_textbox(
        slide, Inches(0.95), y, Inches(11.45), Inches(0.52), text=text,
        font_size=12, color=WHITE, anchor=MSO_ANCHOR.MIDDLE,
    )


def _inv_numbered_rows(slide, items, top: float = 1.7, pitch: float = 0.95):
    """Numbered dataviz circles + lavender pill rows (Next Steps style)."""
    for i, txt in enumerate(items):
        accent = INV_ACCENTS[i % len(INV_ACCENTS)]
        ry = top + i * pitch
        circ = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.8), Inches(ry), Inches(0.53), Inches(0.53)
        )
        circ.fill.solid()
        circ.fill.fore_color.rgb = accent
        circ.line.fill.background()
        _add_textbox(
            slide, Inches(0.8), Inches(ry), Inches(0.53), Inches(0.53), text=str(i + 1),
            font_size=18, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        px, pw, ph = Inches(1.62), Inches(11.1), Inches(0.67)
        py = Inches(ry - 0.07)
        _add_rounded_rect(slide, px, py, pw, ph, INV_PILL)
        _add_filled_rect(slide, px, py, Inches(0.08), ph, accent)
        _add_textbox(
            slide, px + Inches(0.26), py, pw - Inches(0.45), ph, text=txt,
            font_size=14, color=NAVY, anchor=MSO_ANCHOR.MIDDLE,
        )


def _fmt_currency(val, currency=None, compact=False):
    """Format a number as currency. compact=True for $1.2M style.

    ``currency`` defaults to the active plan currency (set by ``generate_pptx``);
    pass an explicit ISO code (e.g. "USD") to force a specific currency for
    intentionally USD-coded benchmark figures.
    """
    if val is None:
        return "—"
    try:
        val = float(val)
    except (TypeError, ValueError):
        return str(val)
    sym = _cur_symbol(currency)
    # display_format.fmt_money never leaves a trailing ".0" ("$150K", never
    # "$150.0K") -- match that here too (fmt_money itself is USD-symbol-only,
    # so it can't be reused directly for a non-USD active-currency plan).
    if compact and abs(val) >= 1_000_000:
        body = f"{val/1_000_000:,.1f}".rstrip("0").rstrip(".")
        return f"{sym}{body}M"
    if compact and abs(val) >= 1_000:
        body = f"{val/1_000:,.1f}".rstrip("0").rstrip(".")
        return f"{sym}{body}K"
    if val == int(val):
        return f"{sym}{int(val):,}"
    return f"{sym}{val:,.2f}"


def _fmt_pct(val, decimals=1):
    """Format as percentage."""
    if val is None:
        return "—"
    try:
        val = float(val)
    except (TypeError, ValueError):
        return str(val)
    if val <= 1.0:  # assume it's a decimal like 0.05 or 1.0 (= 100%)
        val = val * 100
    return f"{val:.{decimals}f}%"


def _fmt_range(low, high, fmt_fn=None):
    """Format a range."""
    fmt = fmt_fn or _fmt_currency
    return f"{fmt(low)} - {fmt(high)}"


def _mark_usd(text: str) -> str:
    """Prefix every bare ``$`` in ``text`` with ``US`` (-> ``US$``).

    Used for fixed US-calibrated benchmark strings (rule b) so the USD
    marker sits DIRECTLY ON the figure itself rather than only in a nearby
    label/caption -- ``US$15 - US$56`` instead of a bare ``$15 - $56``. A
    no-op on values already prefixed with a currency letter (e.g. ``NZ$``).
    """
    if not text:
        return text
    return re.sub(r"(?<![A-Za-z])\$", "US$", text)


def _set_font(
    run,
    size: int = 10,
    bold: bool = False,
    italic: bool = False,
    color: RGBColor = DARK_TEXT,
    name: str = FONT_FAMILY,
):
    """Configure font properties on a text run."""
    run.font.name = name
    # Readability floor: the reference deck never goes below ~9pt. Clamp to 8pt so
    # no caption/label renders as the old unreadable 6-7pt micro-text.
    run.font.size = Pt(max(8, size))
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def _add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str = "",
    font_size: int = 10,
    bold: bool = False,
    italic: bool = False,
    color: RGBColor = DARK_TEXT,
    alignment=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    word_wrap: bool = True,
):
    """Add a text box to a slide and return (shape, text_frame)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None
    try:
        tf.paragraphs[0].alignment = alignment
    except Exception:
        pass
    txBox.text_frame.paragraphs[0].space_before = Pt(0)
    txBox.text_frame.paragraphs[0].space_after = Pt(0)
    try:
        txBox.text_frame._txBody.bodyPr.set(
            "anchor",
            {
                MSO_ANCHOR.TOP: "t",
                MSO_ANCHOR.MIDDLE: "ctr",
                MSO_ANCHOR.BOTTOM: "b",
            }.get(anchor, "t"),
        )
    except Exception:
        pass

    if text:
        p = tf.paragraphs[0]
        p.alignment = alignment
        run = p.add_run()
        run.text = str(text) if text is not None else ""
        _set_font(run, size=font_size, bold=bold, italic=italic, color=color)

    return txBox, tf


# ---------------------------------------------------------------------------
# Measure-then-place autoshrink helpers (O1)
# ---------------------------------------------------------------------------
# python-pptx cannot query a rendering engine for real glyph metrics, so we
# estimate line count from an average character advance keyed to font size.
# For the Poppins/Inter families used here, the mean advance of a bold cap-heavy
# string is ~0.52em and a regular body run ~0.50em; we use 0.53em as a safe
# upper bound so we never UNDER-count lines (which would let text overflow).
# These helpers underpin the exec-summary headline clamp, the KPI band value
# fit, and any single-line value that must not wrap.

_AVG_CHAR_EM = 0.53  # conservative average glyph advance as a fraction of pt size


def _estimate_lines(
    text: str, width_in: float, font_pt: float, char_em: float = _AVG_CHAR_EM
) -> int:
    """Estimate how many wrapped lines ``text`` needs in a box ``width_in`` wide.

    Wrapping is approximated word-by-word using an average character advance;
    it deliberately over-estimates slightly so callers shrink rather than clip.
    ``char_em`` lets callers widen the advance for large, bold, digit-heavy runs
    (e.g. KPI values) whose glyphs are wider than average body prose.
    """
    if not text or width_in <= 0 or font_pt <= 0:
        return 1
    char_w_in = (char_em * font_pt) / 72.0  # pt -> inches
    if char_w_in <= 0:
        return 1
    chars_per_line = max(1, int(width_in / char_w_in))
    lines = 1
    cur = 0
    for word in str(text).split():
        wlen = len(word)
        # account for the joining space when not at line start
        add = wlen + (1 if cur > 0 else 0)
        if cur > 0 and cur + add > chars_per_line:
            lines += 1
            cur = wlen
        else:
            cur += add
        # a single word longer than the line still wraps internally
        while cur > chars_per_line:
            lines += 1
            cur -= chars_per_line
    return max(1, lines)


def _fit_font_to_lines(
    text: str,
    width_in: float,
    start_pt: float,
    max_lines: int,
    min_pt: float = 8.0,
    char_em: float = _AVG_CHAR_EM,
) -> float:
    """Return the largest font size <= ``start_pt`` (>= ``min_pt``, honouring the
    8pt readability floor) at which ``text`` fits within ``max_lines`` in a box
    ``width_in`` wide. Used to clamp headlines / multi-line blocks."""
    size = float(start_pt)
    while size > min_pt and _estimate_lines(text, width_in, size, char_em) > max_lines:
        size -= 0.5
    return max(min_pt, size)


# Big, bold, digit-heavy KPI/hero numerals render noticeably wider than average
# body prose; use a fatter advance so single-line fitting never under-shrinks
# them (which is what let "NZ$29.83" wrap to two lines in the KPI band).
_KPI_CHAR_EM = 0.62


def _fit_font_single_line(
    text: str,
    width_in: float,
    start_pt: float,
    min_pt: float = 8.0,
    char_em: float = _KPI_CHAR_EM,
) -> float:
    """Return the largest font size <= ``start_pt`` at which ``text`` fits on ONE
    line in a box ``width_in`` wide. Keyed to the rendered string's length so a
    long currency-prefixed KPI value (e.g. ``NZ$33.95``) never wraps to 2 lines
    in a single-line band."""
    return _fit_font_to_lines(
        text, width_in, start_pt, max_lines=1, min_pt=min_pt, char_em=char_em
    )


def _autofit_textframe(tf, width_in: float, max_height_in: float, min_pt: float = 8.0):
    """Shrink every run's font (and proportionally its paragraph spacing) until
    the text frame's estimated rendered height fits within ``max_height_in``.

    Used to clamp the exec-summary SCR card bodies (SITUATION / COMPLICATION /
    RESOLUTION) so no row escapes below the card's rounded bottom edge. Estimates
    height as sum over paragraphs of (wrapped-line count * line-height) +
    space_before + space_after, then scales all font sizes by a single factor so
    the layout stays visually uniform. Honors the 8pt readability floor.
    """
    usable_w = max(0.1, width_in - 0.2)  # default 0.1in L/R text insets
    # Conservative rendering constants calibrated against the Keynote render of
    # the exec-summary cards: Poppins/Inter wrap slightly wider than an "average"
    # advance and their single-line box is ~1.4x the point size once leading and
    # top/bottom text insets are included. Under-estimating here is what let the
    # RESOLUTION card's tail rows spill below the card, so we bias toward
    # over-estimating (shrink rather than clip).
    _wrap_em = 0.56  # advance for body prose in these cards
    _line_h_factor = 1.42  # rendered line box height / point size
    _v_inset = 0.10  # text frame top+bottom inset budget (inches)

    def _measure(scale: float) -> float:
        total = _v_inset
        for p in tf.paragraphs:
            runs = [r for r in p.runs if r.text]
            if not runs:
                # empty paragraph still contributes a blank line at base size
                total += (min_pt * _line_h_factor) / 72.0
                continue
            # largest run size in the paragraph drives its line height
            base_pt = max((r.font.size.pt if r.font.size else 10) for r in runs)
            line_pt = base_pt * scale
            ptext = "".join(r.text for r in runs)
            n_lines = _estimate_lines(ptext, usable_w, line_pt, char_em=_wrap_em)
            line_h = (line_pt * _line_h_factor) / 72.0
            sb = (p.space_before.pt if p.space_before else 0) / 72.0
            sa = (p.space_after.pt if p.space_after else 0) / 72.0
            total += n_lines * line_h + sb + sa
        return total

    if _measure(1.0) <= max_height_in:
        return  # already fits at full size

    # Binary-ish search for the largest scale that fits.
    scale = 1.0
    while scale > 0.55 and _measure(scale) > max_height_in:
        scale -= 0.04

    def _apply(scale_val: float) -> None:
        for p in tf.paragraphs:
            for r in p.runs:
                cur = r.font.size.pt if r.font.size else 10
                r.font.size = Pt(max(min_pt, cur * scale_val))
            # scale paragraph spacing too so rows tighten proportionally
            if p.space_before:
                p.space_before = Pt(max(0, p.space_before.pt * scale_val))
            if p.space_after:
                p.space_after = Pt(max(0, p.space_after.pt * scale_val))

    def _measure_current() -> float:
        """Re-measure using each run's CURRENT (already-floored) size."""
        total = _v_inset
        for p in tf.paragraphs:
            runs = [r for r in p.runs if r.text]
            if not runs:
                total += (min_pt * _line_h_factor) / 72.0
                continue
            base_pt = max((r.font.size.pt if r.font.size else 10) for r in runs)
            ptext = "".join(r.text for r in runs)
            n_lines = _estimate_lines(ptext, usable_w, base_pt, char_em=_wrap_em)
            line_h = (base_pt * _line_h_factor) / 72.0
            sb = (p.space_before.pt if p.space_before else 0) / 72.0
            sa = (p.space_after.pt if p.space_after else 0) / 72.0
            total += n_lines * line_h + sb + sa
        return total

    _apply(scale)

    # The per-run 8pt readability floor can prevent the uniform scale from
    # actually reaching ``max_height_in`` when a card is content-overloaded
    # (e.g. the RESOLUTION card's thesis + 5 channels + ML line + goals + cited
    # 2026 block). Rather than let the tail spill below the card, drop trailing
    # paragraphs (lowest-priority content, added last) until it genuinely fits.
    body = tf._txBody
    while _measure_current() > max_height_in and len(tf.paragraphs) > 1:
        last_p = tf.paragraphs[-1]._p
        body.remove(last_p)
    # Don't leave a dangling section header (e.g. "2026 Market Data:") whose
    # content was just trimmed away -- drop it too so the card ends cleanly.
    while len(tf.paragraphs) > 1 and tf.paragraphs[-1].text.strip().endswith(":"):
        body.remove(tf.paragraphs[-1]._p)


def _add_filled_rect(slide, left, top, width, height, fill_color: RGBColor):
    """Add a rectangle shape with a solid fill and no border."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _add_rounded_rect(slide, left, top, width, height, fill_color: RGBColor):
    """Add a rounded rectangle shape with a solid fill and no border."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _add_oval(slide, left, top, width, height, fill_color: RGBColor):
    """Add an oval/circle shape with solid fill and no border."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _add_rule_line(slide, left_inches, top_inches, width_inches, color_hex="202058"):
    """Add a thin horizontal rule line to a slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left_inches),
        Inches(top_inches),
        Inches(width_inches),
        Pt(1.5),  # thin line
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(color_hex)
    shape.line.fill.background()  # no border
    return shape


def _add_paragraph(
    tf,
    text,
    font_size=10,
    bold=False,
    italic=False,
    color=DARK_TEXT,
    alignment=PP_ALIGN.LEFT,
    space_before=0,
    space_after=2,
):
    """Append a paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = str(text) if text is not None else ""
    _set_font(run, size=font_size, bold=bold, italic=italic, color=color)
    return p


def _add_multi_run_paragraph(
    tf, runs_data: List[Tuple], alignment=PP_ALIGN.LEFT, space_before=0, space_after=2
):
    """Add a paragraph with multiple styled runs.
    runs_data: list of (text, font_size, bold, color) tuples.
    """
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    for text, font_size, bold, color in runs_data:
        run = p.add_run()
        run.text = str(text) if text is not None else ""
        _set_font(run, size=font_size, bold=bold, color=color)
    return p


def _kb_extract_range(node: Any) -> str:
    """Pull a display-ready range/figure string out of one
    recruitment_benchmarks.industry_benchmarks[industry][metric] node.

    These KB nodes are free-form dicts (see data loaded via
    kb_loader.load_knowledge_base()): a "range" sub-key when present,
    otherwise the newest-looking numeric-year sub-key (e.g. "2025"), else
    "median"/"total_cost_per_hire"/"recruitment_marketing_only", else the
    first string-valued sub-key found. Returns "" when nothing usable.
    """
    if isinstance(node, str):
        return node
    if not isinstance(node, dict):
        return ""
    for key in ("range", "total_cost_per_hire", "median", "recruitment_marketing_only"):
        val = node.get(key)
        if isinstance(val, str) and val.strip():
            return val.strip()
    year_keys = [k for k in node if isinstance(k, str) and re.fullmatch(r"20\d{2}", k)]
    if year_keys:
        val = node.get(max(year_keys))
        if isinstance(val, str) and val.strip():
            return val.strip()
    for val in node.values():
        if isinstance(val, str) and val.strip():
            return val.strip()
    return ""


def _classify_apply_rate_suffix(apply_rate_str: str) -> str:
    """Return the "(above average - strength)" / "(at industry average)" /
    "(below average - challenge)" style annotation for an apply-rate
    benchmark string, computed strictly from THAT SAME string.

    Two rules, both anchored to the one benchmark being displayed (fixes
    copy:both#1 / visual:atria#1 -- a self-contradictory "lowest of all 24
    occupations, declining ... above average - strength" label):

    1. Only digits immediately followed by "%" count as apply-rate values.
       A naive "grab every digit in the string" parse previously picked up
       the "24" in "(lowest of all 24 occupations, declining)" and averaged
       it with the real 3.22% figure, pushing the average past the
       strength threshold.
    2. If the KB text itself already flags this occupation as the lowest/
       declining (a negative trend baked into the same benchmark data being
       rendered), that verdict wins outright -- it can never be paired with
       a "strength" or "average" label, regardless of the raw percentage.
    """
    if not apply_rate_str:
        return ""
    pct_values = [float(m) for m in re.findall(r"([\d.]+)\s*%", apply_rate_str)]
    avg_rate = sum(pct_values) / len(pct_values) if pct_values else 0.0

    lowered = apply_rate_str.lower()
    negative_trend = any(
        kw in lowered for kw in ("lowest", "declin", "worst", "bottom of")
    )
    if negative_trend:
        return "(below industry average — plan compensates via volume channels)"
    if avg_rate > 5.0:
        return "(above average - strength)"
    if avg_rate >= 2.0:
        return "(at industry average)"
    return "(below average - challenge)"


def _kb_recruitment_industry_benchmark(
    industry: str, data: Optional[Dict]
) -> Optional[Dict[str, str]]:
    """Read cpa/cpc/cph/apply_rate for ``industry`` from
    ``recruitment_benchmarks.industry_benchmarks`` -- the SAME knowledge-base
    section excel_v2's "Recruitment Benchmarks" table reads (see
    excel_v2._build_sheet_executive_summary's ``kb_benchmarks`` lookup).
    Returns ``None`` when the KB or this industry's entry is unavailable.
    """
    kb = (data or {}).get("_knowledge_base") if isinstance(data, dict) else None
    if not kb:
        try:
            from kb_loader import load_knowledge_base

            kb = load_knowledge_base()
        except Exception:  # noqa: BLE001
            return None
    if not isinstance(kb, dict):
        return None
    ind_bm = (
        kb.get("recruitment_benchmarks", {}).get("industry_benchmarks", {})
        if isinstance(kb.get("recruitment_benchmarks"), dict)
        else {}
    )
    entry = ind_bm.get(industry) if isinstance(ind_bm, dict) else None
    if not isinstance(entry, dict):
        return None
    _cpc_node = entry.get("cpc")
    out = {
        "cpa": _kb_extract_range(entry.get("cpa")),
        "cpc": _kb_extract_range(entry.get("cpc")),
        "cph": _kb_extract_range(entry.get("cph")),
        "apply_rate": _kb_extract_range(entry.get("apply_rate")),
        # consistency:atria#4: the CPC YoY trend quoted anywhere in the deck
        # must be this SAME "trend_yoy" figure (not trend_engine's
        # independently-computed live estimate) -- the workbook's benchmark
        # rows quote this exact KB field.
        "cpc_trend_yoy": (
            str(_cpc_node.get("trend_yoy") or "").strip()
            if isinstance(_cpc_node, dict)
            else ""
        ),
    }
    return out if any(out.values()) else None


def _joveo_total_active_publishers(data: Optional[Dict]) -> Optional[int]:
    """Real total-active-publisher count from ``data/joveo_publishers.json``
    (the SAME file ``channel_recommender._enrich_with_joveo_publishers``
    reads) via the shared KB loader, or ``data["_joveo_publishers"]`` when
    the pipeline already attached it.

    strategy:manpower#2: the deck used to fall back to a hardcoded 10238
    default that appears nowhere in the plan data or the KB -- the real
    figure in ``data/joveo_publishers.json`` is 1238. Returns ``None`` (never
    a guessed number) when no real count is available, so the caller can
    drop the statistic rather than invent one.
    """
    pubs = (data or {}).get("_joveo_publishers") if isinstance(data, dict) else None
    if not isinstance(pubs, dict) or not pubs.get("total_active_publishers"):
        kb = (data or {}).get("_knowledge_base") if isinstance(data, dict) else None
        if not isinstance(kb, dict):
            try:
                from kb_loader import load_knowledge_base

                kb = load_knowledge_base()
            except Exception:  # noqa: BLE001
                kb = None
        pubs = kb.get("joveo_publishers") if isinstance(kb, dict) else None
    if not isinstance(pubs, dict):
        return None
    total = pubs.get("total_active_publishers")
    try:
        total = int(total)
    except (TypeError, ValueError):
        return None
    return total if total > 0 else None


def _get_benchmarks(industry: str, data: Optional[Dict] = None) -> Dict[str, str]:
    """Return benchmark data for the given industry.

    v3 resolution cascade:
    1. Synthesized ad_platform_analysis (live API data)
    2. trend_engine.py dynamic benchmarks (with YoY trend arrows)
    3. Hardcoded BENCHMARKS dict (static fallback)

    Returns dict with keys: cpa, cpc, cph, apply_rate, plus optional
    cpc_trend, cpa_trend (e.g. "+8.2% YoY") and confidence.
    """
    # Start with hardcoded fallback
    result = dict(BENCHMARKS.get(industry, BENCHMARKS["general_entry_level"]))
    result["confidence"] = "curated"
    # S3: track whether cpa/cpc/cph currently hold a fixed US-benchmark
    # constant (rule b -- needs an inline "(USD)" marker at render time) vs a
    # plan-derived figure (rule a -- localize to the active plan currency).
    # BENCHMARKS/appcast/trend_engine are static US-calibrated lookup tables;
    # ad_platform_analysis (live_api) and the budget-engine CPH override below
    # are the PLAN'S OWN projected figures.
    result["cpa_is_usd_benchmark"] = True
    result["cpc_is_usd_benchmark"] = True
    result["cph_is_usd_benchmark"] = True

    # Layer 2: trend_engine dynamic benchmarks (v3)
    if _HAS_TREND_ENGINE:
        try:
            campaign_month = 0
            if data:
                try:
                    campaign_month = int(data.get("campaign_start_month") or 0 or 0)
                except (ValueError, TypeError):
                    campaign_month = 0
            current_month = (
                campaign_month
                if (campaign_month and 1 <= campaign_month <= 12)
                else datetime.datetime.now().month
            )
            # Get CPC across platforms
            cpc_vals = []
            cpa_vals = []
            trend_dirs = []
            trend_yoys = []
            for plat in ["google", "meta_fb", "indeed", "linkedin", "programmatic"]:
                cpc_result = _trend_engine.get_benchmark(
                    platform=plat,
                    industry=industry or "general",
                    metric="cpc",
                    collar_type="both",
                    location="",
                    month=current_month,
                )
                if cpc_result and isinstance(cpc_result, dict):
                    v = cpc_result.get("value") or 0
                    if isinstance(v, (int, float)) and v > 0:
                        cpc_vals.append(v)
                        td = cpc_result.get("trend_direction", "stable")
                        ty = cpc_result.get("trend_pct_yoy") or 0
                        trend_dirs.append(td)
                        trend_yoys.append(ty)
                cpa_result = _trend_engine.get_benchmark(
                    platform=plat,
                    industry=industry or "general",
                    metric="cpa",
                    collar_type="both",
                    location="",
                    month=current_month,
                )
                if cpa_result and isinstance(cpa_result, dict):
                    v = cpa_result.get("value") or 0
                    if isinstance(v, (int, float)) and v > 0:
                        cpa_vals.append(v)

            if cpc_vals:
                min_cpc = min(cpc_vals)
                max_cpc = max(cpc_vals)
                result["cpc"] = (
                    f"${min_cpc:.2f} - ${max_cpc:.2f}"
                    if len(cpc_vals) > 1
                    else f"${min_cpc:.2f}"
                )
                result["confidence"] = "trend_engine"
                # Trend annotation
                if trend_yoys:
                    avg_yoy = sum(trend_yoys) / len(trend_yoys)
                    arrow = "+" if avg_yoy > 0 else ""
                    dom_trend = (
                        max(set(trend_dirs), key=trend_dirs.count)
                        if trend_dirs
                        else "stable"
                    )
                    result["cpc_trend"] = f"{arrow}{avg_yoy:.1f}% YoY"
                    result["cpc_trend_direction"] = dom_trend
            if cpa_vals:
                min_cpa = min(cpa_vals)
                max_cpa = max(cpa_vals)
                result["cpa"] = (
                    f"${min_cpa:.0f} - ${max_cpa:.0f}"
                    if len(cpa_vals) > 1
                    else f"${min_cpa:.0f}"
                )
        except Exception:
            pass  # Fall through to synthesized or static

    # Layer 1.5: recruitment_benchmarks.industry_benchmarks KB overlay --
    # single-sources the "Industry CPA/CPC/Apply Rate" figures with the SAME
    # KB section excel_v2._build_sheet_executive_summary reads for its
    # "Recruitment Benchmarks" table (kb['recruitment_benchmarks']
    # ['industry_benchmarks'][industry]). Without this, the deck's numbers
    # came from trend_engine's independent static tables while the workbook
    # quoted this KB section, so the two artifacts could show different CPA/
    # CPC ranges for the identical industry. Runs AFTER trend_engine (so it
    # overrides those numbers when this KB section exists for the industry)
    # but BEFORE the live ad_platform_analysis / budget-engine layers below
    # (so the plan's own live figures, when present, still win over any
    # benchmark table).
    if data:
        try:
            _mi_bm = _kb_recruitment_industry_benchmark(industry, data)
        except Exception:  # noqa: BLE001 -- never break benchmark resolution
            _mi_bm = None
        if _mi_bm:
            if _mi_bm.get("cpa"):
                result["cpa"] = _mi_bm["cpa"]
            if _mi_bm.get("cpc"):
                result["cpc"] = _mi_bm["cpc"]
            if _mi_bm.get("cph"):
                result["cph"] = _mi_bm["cph"]
            if _mi_bm.get("apply_rate"):
                result["apply_rate"] = _mi_bm["apply_rate"]
            if _mi_bm.get("cpc_trend_yoy"):
                # consistency:atria#4: override trend_engine's independent
                # "+9.1% YoY" live estimate with the KB's own "+10-15%"
                # trend_yoy figure -- the same one the workbook quotes --
                # so the deck never contradicts the range the workbook's
                # benchmark rows carry.
                result["cpc_trend"] = _mi_bm["cpc_trend_yoy"]
                if result["cpc_trend"].strip().startswith("-"):
                    result["cpc_trend_direction"] = "falling"
                elif result["cpc_trend"].strip().startswith("+"):
                    result["cpc_trend_direction"] = "rising"
                else:
                    result["cpc_trend_direction"] = "stable"
            else:
                # strategy:manpower#2: this industry HAS a real KB benchmark
                # table (recruitment_benchmarks_deep.json), but its cpc node
                # carries no "trend_yoy" field -- e.g. logistics_supply_chain,
                # whose workbook CPC benchmark has no YoY trend at all. Never
                # let trend_engine's independently-computed live estimate
                # (e.g. a "+9.5% YoY" that appears nowhere in the workbook or
                # KB) leak into a client-facing deck as if it were sourced.
                # Drop the untraceable number rather than show one that
                # can't be traced to any KB/workbook cell.
                result.pop("cpc_trend", None)
                result.pop("cpc_trend_direction", None)
            result["confidence"] = "market_intelligence_kb"

    # Layer 1: Synthesized ad_platform_analysis overrides (live API data)
    if data:
        synthesized = data.get("_synthesized", {})
        if isinstance(synthesized, dict):
            ad_plat = synthesized.get("ad_platform_analysis", {})
            if isinstance(ad_plat, dict) and ad_plat:
                live_cpcs = []
                live_cpas = []
                for plat_name, plat_data in ad_plat.items():
                    if not isinstance(plat_data, dict) or plat_name.startswith("_"):
                        continue
                    cpc = plat_data.get("avg_cpc") or plat_data.get("cpc")
                    cpa = plat_data.get("avg_cpa") or plat_data.get("cpa")
                    if cpc and isinstance(cpc, (int, float)) and cpc > 0:
                        live_cpcs.append(cpc)
                    if cpa and isinstance(cpa, (int, float)) and cpa > 0:
                        live_cpas.append(cpa)
                if live_cpcs:
                    min_cpc = min(live_cpcs)
                    max_cpc = max(live_cpcs)
                    # S3: plan's own live ad-platform CPC projection -- localize.
                    result["cpc"] = (
                        f"{_fmt_currency(min_cpc)} - {_fmt_currency(max_cpc)}"
                        if min_cpc != max_cpc
                        else _fmt_currency(min_cpc)
                    )
                    result["confidence"] = "live_api"
                    result["cpc_is_usd_benchmark"] = False
                if live_cpas:
                    min_cpa = min(live_cpas)
                    max_cpa = max(live_cpas)
                    # S3: plan's own live ad-platform CPA projection -- localize.
                    result["cpa"] = (
                        f"{_fmt_currency(round(min_cpa))} - {_fmt_currency(round(max_cpa))}"
                        if min_cpa != max_cpa
                        else _fmt_currency(round(min_cpa))
                    )
                    result["confidence"] = "live_api"
                    result["cpa_is_usd_benchmark"] = False

    # Layer 0: Budget engine CPH and apply_rate overrides
    # These were NEVER overridden before (always hardcoded) -- fix v3.4.1
    if data:
        budget_alloc = data.get("_budget_allocation", {})
        if isinstance(budget_alloc, dict) and budget_alloc:
            total_proj = budget_alloc.get("total_projected", {})
            if isinstance(total_proj, dict):
                # CPH from budget engine
                live_cph = total_proj.get("cost_per_hire") or total_proj.get("cph")
                if isinstance(live_cph, (int, float)) and live_cph > 0:
                    # Format as range: computed +/- 20% to show realistic spread
                    # S3: this is the plan's OWN cost-per-hire -- localize.
                    low_cph = live_cph * 0.8
                    high_cph = live_cph * 1.2
                    result["cph"] = (
                        f"{_fmt_currency(low_cph)} - {_fmt_currency(high_cph)}"
                    )
                    result["cph_is_usd_benchmark"] = False
                # Apply rate from budget engine
                live_apply_rate = total_proj.get("apply_rate") or total_proj.get(
                    "conversion_rate"
                )
                if isinstance(live_apply_rate, (int, float)) and live_apply_rate > 0:
                    # Apply rate might be 0-1 or 0-100; normalize
                    if live_apply_rate < 1:
                        live_apply_rate *= 100
                    low_ar = live_apply_rate * 0.85
                    high_ar = live_apply_rate * 1.15
                    result["apply_rate"] = f"{low_ar:.1f}% - {high_ar:.1f}%"

        # Also check synthesized salary_intelligence for CPH context
        # but ONLY if the budget engine didn't already set it (budget engine
        # computes CPH from actual budget/projections -- higher quality).
        synthesized = data.get("_synthesized", {})
        _budget_set_cph = (
            isinstance(budget_alloc, dict) and budget_alloc
        )  # budget engine ran
        if not _budget_set_cph and isinstance(synthesized, dict):
            salary_intel = synthesized.get("salary_intelligence", {})
            if isinstance(salary_intel, dict):
                synth_cph = salary_intel.get(
                    "estimated_cost_per_hire"
                ) or salary_intel.get("cph")
                if isinstance(synth_cph, (int, float)) and synth_cph > 0:
                    # S3: derived from the plan's own salary data -- localize.
                    low_cph = synth_cph * 0.8
                    high_cph = synth_cph * 1.2
                    result["cph"] = (
                        f"{_fmt_currency(low_cph)} - {_fmt_currency(high_cph)}"
                    )
                    result["cph_is_usd_benchmark"] = False

    # Layer 2.5: Appcast 2026 KB Benchmark Overlay (Priority 3)
    # Cross-reference/override CPA, CPH, apply_rate with Appcast occupation-level
    # data if the current values are still from hardcoded fallback.
    if result.get("confidence") == "curated":
        try:
            _PPT_APPCAST_MAP = {
                "healthcare": "healthcare",
                "healthcare_medical": "healthcare",
                "technology": "technology",
                "tech_engineering": "technology",
                "retail": "retail",
                "retail_consumer": "retail",
                "finance": "finance",
                "finance_banking": "finance",
                "insurance": "insurance",
                "logistics": "warehousing_logistics",
                "logistics_supply_chain": "warehousing_logistics",
                "transportation": "transportation",
                "hospitality": "hospitality",
                "hospitality_travel": "hospitality",
                "food_beverage": "food_service",
                "blue_collar_trades": "construction_skilled_trades",
                "construction": "construction_skilled_trades",
                "construction_real_estate": "construction_skilled_trades",
                "pharma_biotech": "science_engineering",
                "energy_utilities": "science_engineering",
                "general": "administration",
                "general_entry_level": "administration",
                "government_utilities": "administration",
                "manufacturing": "manufacturing",
                "education": "education",
                "legal_services": "legal",
                "marketing": "marketing_advertising",
                "media_entertainment": "marketing_advertising",
                "professional_services": "consulting",
            }
            _app_key = _PPT_APPCAST_MAP.get(industry) or ""
            if _app_key and data:
                _kb = data.get("_knowledge_base", {})
                if not _kb:
                    try:
                        from kb_loader import load_knowledge_base

                        _kb = load_knowledge_base()
                    except Exception:
                        _kb = {}
                _wp = _kb.get("white_papers", {})
                _appcast_bm = (
                    _wp.get("reports", {})
                    .get("appcast_benchmark_2026", {})
                    .get("benchmarks", {})
                )
                if _appcast_bm:
                    _occ_cpa = _appcast_bm.get("cpa_by_occupation_2025", {}).get(
                        _app_key
                    )
                    _occ_cph = _appcast_bm.get("cph_by_occupation_2025", {}).get(
                        _app_key
                    )
                    _occ_ar = _appcast_bm.get("apply_rate_by_occupation_2025", {}).get(
                        _app_key
                    )
                    if _occ_cpa:
                        result["cpa"] = _occ_cpa
                        result["cpa_source"] = "Appcast 2026 (302M clicks)"
                    if _occ_cph:
                        result["cph"] = _occ_cph
                        result["cph_source"] = "Appcast 2026 (302M clicks)"
                    if _occ_ar:
                        result["apply_rate"] = _occ_ar
                        result["apply_rate_source"] = "Appcast 2026 (302M clicks)"
                    if any([_occ_cpa, _occ_cph, _occ_ar]):
                        result["confidence"] = "appcast_kb"
        except Exception:
            pass  # Non-fatal; fall through to curated values

    return result


def _get_complications(industry: str) -> List[str]:
    """Return complication bullets for the industry, with apply rate framed correctly."""
    base = COMPLICATIONS.get(
        industry,
        [
            "Talent acquisition costs rising across sectors",
            "Competition for qualified candidates intensifying",
            "Traditional sourcing channels showing diminishing returns",
            "Time-to-fill expanding, impacting operational capacity",
        ],
    )

    # Get the apply rate for this industry and frame appropriately
    benchmarks = BENCHMARKS.get(industry, BENCHMARKS.get("general_entry_level", {}))
    apply_rate_str = benchmarks.get("apply_rate") or ""
    if apply_rate_str:
        # Parse apply rate - handle ranges like "3.2% - 4.5%" or single values like "6.41%"
        import re as _re

        rates = _re.findall(r"[\d.]+", apply_rate_str)
        if rates:
            avg_rate = sum(float(r) for r in rates) / len(rates)
            # Only add apply rate as complication if it's genuinely low (below 2%)
            if avg_rate < 2.0:
                base = list(base)  # make mutable copy
                base.append(
                    f"Low {apply_rate_str} apply rate indicates competitive market pressure"
                )

    return base


def _get_industry_comparison(
    industry: str, data: Optional[Dict] = None
) -> Dict[str, Any]:
    """Return industry benchmark comparison data.

    v3.4.1: overlay dynamic data from budget engine / synthesized data
    on top of the hardcoded base, so live API values take precedence.
    """
    result = dict(
        INDUSTRY_BENCHMARKS_COMPARISON.get(
            industry, INDUSTRY_BENCHMARKS_COMPARISON["general_entry_level"]
        )
    )

    if data:
        budget_alloc = data.get("_budget_allocation", {})
        if isinstance(budget_alloc, dict) and budget_alloc:
            total_proj = budget_alloc.get("total_projected", {})
            if isinstance(total_proj, dict):
                live_cpa = total_proj.get("cpa") or total_proj.get(
                    "cost_per_application"
                )
                if isinstance(live_cpa, (int, float)) and live_cpa > 0:
                    # S3: the plan's OWN CPA -- localize.
                    result["avg_cpa"] = (
                        f"{_fmt_currency(live_cpa * 0.8)} - {_fmt_currency(live_cpa * 1.2)}"
                    )
                live_cph = total_proj.get("cost_per_hire") or total_proj.get("cph")
                if isinstance(live_cph, (int, float)) and live_cph > 0:
                    # S3: the plan's OWN cost-per-hire -- localize.
                    result["avg_cph"] = (
                        f"{_fmt_currency(live_cph * 0.8)} - {_fmt_currency(live_cph * 1.2)}"
                    )
            ch_allocs = budget_alloc.get("channel_allocations", {})
            if isinstance(ch_allocs, dict) and ch_allocs:
                n_ch = len(
                    [
                        c
                        for c in ch_allocs.values()
                        if isinstance(c, dict) and (c.get("percentage") or 0) > 0
                    ]
                )
                if n_ch > 0:
                    # Keep as int -- used in arithmetic comparisons downstream
                    result["avg_channels"] = n_ch

    return result


def _is_us_only_campaign(data: Dict) -> bool:
    """Check if all campaign locations are within the United States.

    Thin delegate to plan_geo.is_us_plan -- the single source of truth for
    US/non-US plan classification (see plan_geo.py module docstring for the
    two shipped bugs this consolidation fixes: hard-return-False-on-first-
    unresolvable-candidate, and no concept of bare US state names)."""
    return _plan_geo.is_us_plan(data)


def _selected_channels(data: Dict) -> Dict[str, Dict[str, Any]]:
    """Return only the channels the user toggled on, with redistributed percentages.
    Uses industry-aware allocation profiles for differentiated budget splits.
    Automatically excludes APAC/EMEA channels for US-only campaigns."""
    cats = data.get("channel_categories", {})
    if isinstance(cats, list):
        cats = {k: True for k in cats}

    # Skip international channels for US-only campaigns
    us_only = _is_us_only_campaign(data)
    if us_only:
        cats["apac_regional"] = False
        cats["emea_regional"] = False

    # Get industry-aware base allocation
    industry = data.get("industry", "general_entry_level")
    budget_str = data.get("budget") or ""
    roles = data.get("roles") or []
    num_roles = len(roles) if roles else 0
    locations = data.get("locations") or []
    alloc_base = _get_industry_alloc(
        industry, budget_str, num_roles, roles, locations=locations
    )

    selected = {}
    for key, meta in alloc_base.items():
        if cats.get(key, False):
            selected[key] = dict(meta)

    if not selected:
        # Legacy plans without channel_categories: include ALL channels
        # (New UI defaults channels to OFF, so empty cats means legacy data)
        has_explicit_cats = any(v for v in cats.values()) if cats else False
        if not has_explicit_cats:
            for key, meta in alloc_base.items():
                selected[key] = dict(meta)
            # S48 FIX: Apply geo-filter AFTER legacy fallback populates all channels.
            # When channel_categories is empty (legacy/API path), the us_only check
            # at the top correctly sets cats["apac_regional"] = False, but the empty
            # cats dict triggers this fallback which re-includes ALL channels.
            # We must strip APAC/EMEA here for US-only campaigns.
            if us_only:
                selected.pop("apac_regional", None)
                selected.pop("emea_regional", None)
        else:
            # cats was provided but no True values matched alloc_base keys;
            # fall back to programmatic + global + social as minimum
            for key in ("programmatic_dsp", "global_boards", "social_media"):
                if key in alloc_base:
                    selected[key] = dict(alloc_base[key])

    raw_total = sum(v["pct"] for v in selected.values())
    if raw_total > 0:
        for v in selected.values():
            v["pct"] = round(v["pct"] / raw_total * 100)
        diff = 100 - sum(v["pct"] for v in selected.values())
        if diff != 0:
            first_key = next(iter(selected))
            selected[first_key]["pct"] += diff

    return selected


def _largest_remainder_round(values: Dict[str, float], total: int = 100) -> Dict[str, int]:
    """Round a dict of {key: float_percentage} to integers that sum EXACTLY to
    ``total`` (default 100), using the largest-remainder (Hamilton) method.

    S6 fix: independently rounding each channel's percentage (``round(pct)``)
    can drift the displayed total to 99% or 101% once 3+ channels are involved
    -- confirmed on the Pratt & Whitney NZ deck where 8 channels' rounded
    percentages summed to 101% on both the slide-5 channel-mix chart and the
    slide-6 breakdown table's Total row. Largest-remainder rounding floors every
    value, then hands the leftover units (total - sum-of-floors) to the entries
    with the largest fractional remainder, one unit each -- guaranteeing an
    exact reconciliation instead of independent per-item drift.

    Non-positive/zero inputs are floored to 0 and excluded from the remainder
    distribution (a zero-share channel should never be bumped to 1% just to
    make the total foot).
    """
    if not values:
        return {}

    # Normalize to `total` first. This function is designed to correct genuine
    # ROUNDING drift, where the raw values already sum very close to `total`
    # (e.g. 99.8 from float percentages) -- in that case the scale factor is
    # ~1.0 and this is a no-op. But if the caller only has a SUBSET of the full
    # channel mix (e.g. two channels were dropped upstream by an unrelated
    # channel-selection bug), the raw values can sum well short of `total`
    # (e.g. 93); scaling first means largest-remainder only ever has to correct
    # a <1-point-per-item rounding remainder, never asked to invent several
    # whole extra points by capping every item at +1 (which would silently
    # under-shoot when leftover > len(values), e.g. 93 -> 99 instead of 100).
    raw_total = sum(max(0.0, float(v or 0)) for v in values.values())
    scale = (total / raw_total) if raw_total > 0 else 1.0

    floors: Dict[str, int] = {}
    remainders: Dict[str, float] = {}
    for k, v in values.items():
        v = max(0.0, float(v or 0)) * scale
        f = math.floor(v)
        floors[k] = int(f)
        remainders[k] = v - f

    assigned = sum(floors.values())
    leftover = int(round(total - assigned))
    if leftover > 0:
        # Largest remainder first; among eligible (non-zero-value) keys only.
        order = sorted(
            (k for k in values if values[k]),
            key=lambda k: remainders[k],
            reverse=True,
        )
        # Guard: with the pre-scale step above, leftover should never exceed
        # the number of eligible keys (each item's remainder is < 1 after
        # scaling to `total`) -- but keep the cap defensive against float
        # edge cases so we never try to bump the same key past +1.
        for k in order[:leftover]:
            floors[k] += 1
    elif leftover < 0:
        # Values summed above `total` even after flooring (shouldn't normally
        # happen, but guard against float-precision surprises) -- claw back
        # from the smallest remainders among entries that still have a
        # positive rounded share, so nothing goes negative.
        order = sorted(
            (k for k in values if floors[k] > 0),
            key=lambda k: remainders[k],
        )
        for k in order[: (-leftover)]:
            floors[k] -= 1

    return floors


def _reconcile_channel_percentages(
    channels: Dict[str, Dict[str, Any]], ba_channel_alloc: Dict[str, Any]
) -> Dict[str, int]:
    """Match each display channel to its budget-engine allocation and return
    largest-remainder-rounded integer percentages keyed by the SAME keys as
    ``channels``, guaranteed to sum to 100 across all of ``channels``.

    Every slide that displays per-channel allocation percentages (slide 5's
    channel-mix bar chart and category-attribution band, slide 6's
    channel-by-channel table and embedded pie chart) should call this ONCE
    per builder instead of rounding ``ba_match.get("percentage")`` independently
    -- that per-call independent rounding is what let the displayed total drift
    to 101% on an 8-channel plan even though the underlying dollar amounts
    summed to exactly the stated budget.
    """
    if not isinstance(ba_channel_alloc, dict) or not ba_channel_alloc:
        return {}
    raw_pcts: Dict[str, float] = {}
    for ch_key, ch_data in channels.items():
        ba_match = ba_channel_alloc.get(ch_key)
        if not ba_match:
            ch_label_lower = (ch_data.get("label") or "").lower()
            for ba_key, ba_val in ba_channel_alloc.items():
                if isinstance(ba_val, dict):
                    ba_label = ba_val.get("label", ba_key).lower()
                    if ba_label == ch_label_lower or ba_key.lower() == ch_key.lower():
                        ba_match = ba_val
                        break
        if ba_match and isinstance(ba_match, dict):
            real_pct = ba_match.get("percentage") or 0
            if real_pct > 0:
                raw_pcts[ch_key] = real_pct
    if not raw_pcts:
        return {}
    return _largest_remainder_round(raw_pcts)


def _goal_labels(data: Dict) -> List[str]:
    """Return human-readable campaign goal labels."""
    goals = data.get("campaign_goals") or []
    return [GOAL_LABELS.get(g, g.replace("_", " ").title()) for g in goals]


def _parse_budget_number(budget_str) -> Optional[float]:
    """Try to extract a numeric budget value from a string.

    Delegates to shared_utils.parse_budget_display for consistent parsing
    across all modules.
    """
    return parse_budget_display(budget_str)


def _format_budget_display(budget_str: str) -> str:
    """Format budget for hero stat display (uses the active plan currency)."""
    val = _parse_budget_number(budget_str)
    if val is None:
        return budget_str
    sym = _cur_symbol()
    if val >= 1000000:
        body = f"{val / 1000000:.1f}".rstrip("0").rstrip(".")
        return f"{sym}{body}M"
    if val >= 1000:
        return f"{sym}{val / 1000:.0f}K"
    return f"{sym}{val:,.0f}"


def _channel_categories_grouped(channels: Dict) -> Dict[str, List[Dict]]:
    """Group channels by their category for attribution diagram."""
    groups: Dict[str, List[Dict]] = {}
    for key, ch in channels.items():
        cat = ch.get("category", "Other")
        if cat not in groups:
            groups[cat] = []
        groups[cat].append(ch)
    return groups


def _add_footer(slide, today: str):
    """Add the single canonical Joveo footer band (date left, credit right).

    This is the ONLY footer drawn on a slide. ``_add_data_sources_footnote`` and
    ``_add_enrichment_badge`` intentionally do NOT draw their own footer/credit,
    so the bottom margin never stacks overlapping captions (the previous version
    drew three separate 6-7pt lines at y~7.0 that collided on every slide).
    """
    rule_y = Inches(7.12)
    # Thin warm-gray rule (matches the reference deck's minimal footer treatment)
    _add_filled_rect(slide, Inches(0.55), rule_y, Inches(12.23), Inches(0.012), WARM_GRAY)
    # Left-aligned date -- 9pt readability floor (was 7pt)
    _add_textbox(
        slide,
        Inches(0.55),
        rule_y + Inches(0.05),
        Inches(4.5),
        Inches(0.26),
        text=today,
        font_size=9,
        color=MUTED_TEXT,
        alignment=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    # Right-aligned single attribution line in Downy Teal
    _add_textbox(
        slide,
        Inches(5.75),
        rule_y + Inches(0.05),
        Inches(7.03),
        Inches(0.26),
        text="Created by Shubham Singh Chandel  •  Powered by Joveo's Global Supply Team",
        font_size=9,
        color=TEAL,
        alignment=PP_ALIGN.RIGHT,
        anchor=MSO_ANCHOR.MIDDLE,
    )


_TITLE_ACRONYMS = {
    "Ai": "AI", "Roi": "ROI", "Cpa": "CPA", "Cpc": "CPC", "Cph": "CPH",
    "Dsp": "DSP", "Crm": "CRM", "Kpi": "KPI", "Us": "US", "Uk": "UK",
    "Eu": "EU", "Seo": "SEO", "Ats": "ATS", "Ppc": "PPC", "Qc": "QC",
    "Roas": "ROAS", "Dei": "DEI",
}


def _smart_title(s: str) -> str:
    """Title-case a heading while preserving known acronyms (AI, ROI, CPA...)."""
    return " ".join(_TITLE_ACRONYMS.get(w, w) for w in str(s).title().split())


def _add_top_band(slide, left_text: str, right_text: str = "", band_color=NAVY):
    """Invisible-style light header (replaces the old navy nav band).

    Paints the light lavender canvas, the top purple accent rule, the ``joveo``
    wordmark, and a big indigo section title (``left_text``, smart-title-cased).
    ``right_text`` (date/client) is intentionally NOT shown here -- the footer
    carries the date, matching the reference deck. Returns a body-top hint (EMU).
    """
    _inv_canvas(slide, wordmark=False)
    # Compact wordmark + big title sit ABOVE y=0.92 so each builder's existing
    # action/insight line at y=0.92 reads as the subtitle (no per-builder edits,
    # no collision). Section names are short -> one line at 26pt.
    _add_textbox(
        slide,
        Inches(0.63),
        Inches(0.15),
        Inches(4.0),
        Inches(0.3),
        text="joveo",
        font_size=15,
        bold=True,
        color=BLUE,
    )
    _add_textbox(
        slide,
        Inches(0.63),
        Inches(0.45),
        Inches(12.05),
        Inches(0.46),
        text=_smart_title(left_text),
        font_size=26,
        bold=True,
        color=NAVY,
    )
    return Inches(1.45)


def _confidence_color(confidence: str) -> Tuple[RGBColor, str]:
    """Return (color, label) for a confidence level string.

    Green bold for live/high confidence, blue for curated, amber for fallback.
    """
    conf_lower = str(confidence).lower()
    if conf_lower in ("live_api", "high", "synthesized"):
        return (GREEN, "Live Data")
    elif conf_lower in ("trend_engine", "cached_api"):
        return (BLUE, "Trend Engine")
    elif conf_lower in ("curated", "knowledge_base", "medium"):
        return (MEDIUM_BLUE, "Curated")
    else:
        return (AMBER, "Estimated")


def _add_enrichment_badge(slide, enriched):
    """No-op (retained for call-site compatibility).

    The old bottom-right "Live data: ..." badge sat at y=7.1 and overlapped the
    footer credit 95% on every slide that used it. Data provenance now lives in
    the workbook's "Sources & Confidence" sheet and the per-slide confidence
    caption, so the deck footer stays clean and uncluttered.
    """
    return


def _add_data_sources_footnote(slide, data: Dict, benchmarks: Dict):
    """Confidence indicator + optional disclaimers, placed ABOVE the footer rule.

    Does NOT draw a footer/rule/credit — ``_add_footer`` owns the footer band.
    Everything stays at/above y=6.84 so it never collides with the footer or
    runs off the bottom edge (the previous version stacked four 6pt italic lines
    at y=7.0-7.3 that overlapped each other and the footer credit).
    """
    conf = benchmarks.get("confidence", "curated")
    conf_color, conf_label = _confidence_color(conf)

    caption_y = Inches(6.84)
    # Confidence indicator (left)
    _add_textbox(
        slide,
        Inches(0.55),
        caption_y,
        Inches(3.2),
        Inches(0.22),
        text=conf_label,
        font_size=9,
        bold=True,
        color=conf_color,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # Combined disclaimers (right), only when present, truncated to fit one line.
    # NOTE: the "Location advisory: ..." path was removed from the deck
    # entirely (owner call) -- the workbook's Sources & Confidence sheet
    # already carries the full per-location warning; duplicating a truncated
    # version here just repeated it with less detail.
    disclaimers = []
    if not _is_us_only_campaign(data):
        disclaimers.append(
            "Benchmarks are US-calibrated — international markets may vary."
        )
    if disclaimers:
        _add_textbox(
            slide,
            Inches(3.9),
            caption_y,
            Inches(8.85),
            Inches(0.22),
            text=_trunc_word("   ".join(disclaimers), 150),
            font_size=8,
            italic=True,
            color=AMBER,
            alignment=PP_ALIGN.RIGHT,
            anchor=MSO_ANCHOR.MIDDLE,
        )


def _format_salary(amount):
    """Format a salary number into human-readable string like $85K or $125K.

    Uses the active plan currency symbol so non-USD plans render correctly.
    """
    if not isinstance(amount, (int, float)) or amount <= 0:
        return ""
    sym = _cur_symbol()
    if amount >= 1000:
        return f"{sym}{amount / 1000:.0f}K"
    return f"{sym}{amount:,.0f}"


# ===================================================================
# SLIDE 1 - Cover / Section Divider: Title Slide
# ===================================================================


def _build_slide_cover(prs: Presentation, data: Dict):
    """Build a premium full-bleed cover slide with LinkedIn-style section divider pattern."""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    client = data.get("client_name", "Client")
    industry_label = data.get("industry_label") or ""
    today = datetime.date.today().strftime("%B %d, %Y")

    # Full dark indigo background (matches the Invisible title slide)
    _add_filled_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, NAVY)

    # Top-right corner decoration: layered purple-light + teal rectangles
    _add_filled_rect(slide, Inches(10.0), Inches(0), Inches(3.33), Inches(0.2), JOVEO_LIGHT_PURPLE)
    _add_filled_rect(slide, Inches(11.33), Inches(0.2), Inches(2.0), Inches(0.13), TEAL)

    # joveo wordmark, top-left
    _add_textbox(
        slide, Inches(0.62), Inches(0.5), Inches(4), Inches(0.55),
        text="joveo", font_size=26, bold=True, color=JOVEO_LIGHT_PURPLE,
    )

    # Main title
    _add_textbox(
        slide, Inches(0.62), Inches(1.95), Inches(11.5), Inches(0.95),
        text="Media Plan", font_size=52, bold=True, color=WHITE,
    )
    # "for" connector + client hero
    _add_textbox(
        slide, Inches(0.66), Inches(3.0), Inches(4), Inches(0.38),
        text="for", font_size=20, italic=True, color=WHITE,
    )
    _add_textbox(
        slide, Inches(0.62), Inches(3.48), Inches(11.8), Inches(1.0),
        text=client, font_size=42, bold=True, color=LIGHT_TEAL,
    )

    # Industry subtitle
    if industry_label:
        _add_textbox(
            slide, Inches(0.64), Inches(4.5), Inches(10), Inches(0.5),
            text=industry_label, font_size=18, color=LIGHT_TEAL,
        )

    # Company tagline from enrichment data (Wikipedia description)
    enriched = data.get("_enriched", {})
    company_info = enriched.get("company_info", {}) if enriched else {}
    if company_info and company_info.get("description"):
        desc = company_info["description"]
        first_sentence_end = desc.find(".")
        if 0 < first_sentence_end < 120:
            tagline = desc[: first_sentence_end + 1]
        else:
            tagline = desc[:120].rsplit(" ", 1)[0] + "..." if len(desc) > 120 else desc
        _add_textbox(
            slide, Inches(0.64), Inches(5.02), Inches(9), Inches(0.4),
            text=tagline, font_size=11, italic=True, color=LIGHT_MUTED,
        )

    # Prepared-by + date (bottom-left)
    _add_textbox(
        slide, Inches(0.62), Inches(5.62), Inches(8), Inches(0.4),
        text="Created by Shubham Singh Chandel  •  Powered by Joveo's Global Supply Team",
        font_size=13, color=TEAL,
    )
    _add_textbox(
        slide, Inches(0.62), Inches(6.04), Inches(6), Inches(0.4),
        text=today, font_size=13, color=LIGHT_MUTED,
    )

    # Bottom-right tagline
    _add_textbox(
        slide, Inches(8.4), Inches(6.4), Inches(4.3), Inches(0.6),
        text="High-Performance, AI-Led\nRecruitment Marketing Platform",
        font_size=11, color=LIGHT_MUTED, alignment=PP_ALIGN.RIGHT,
    )

    # Bottom edge brand bars (purple-light + teal)
    _add_filled_rect(slide, Inches(0), SLIDE_HEIGHT - Inches(0.16), SLIDE_WIDTH, Inches(0.08), JOVEO_LIGHT_PURPLE)
    _add_filled_rect(slide, Inches(0), SLIDE_HEIGHT - Inches(0.08), SLIDE_WIDTH, Inches(0.08), TEAL)


# ===================================================================
# SLIDE 2 - Executive Summary with Hero Stat
# ===================================================================


def _cap_roles_for_headline(roles: List[Any], cap: int = 2) -> str:
    """Cap a role list for headline prose: at most ``cap`` named roles, then
    'and N more roles' -- never the truncated-grammar 'Role1, Role2, Role3'
    with no indication more roles exist, and never a dangling trailing
    comma."""
    names = [str(r) for r in (roles or []) if str(r).strip()]
    if not names:
        return "key roles"
    if len(names) <= cap:
        if len(names) == 1:
            return names[0]
        return ", ".join(names[:-1]) + " and " + names[-1]
    shown = ", ".join(names[:cap])
    remaining = len(names) - cap
    return f"{shown}, and {remaining} more role{'s' if remaining != 1 else ''}"


def _compute_blended_cph(budget_alloc: Optional[Dict]) -> "tuple[float, int]":
    """Compute the plan's blended cost-per-hire ONCE, from the channel-level
    sums in ``data['_budget_allocation']`` -- the same source-of-truth
    calculation the S48 fix established for the Executive Summary slide's
    hero stat (per-channel projected_hires summed, then total_budget /
    total_hires). Every slide that renders a "blended"/"average" cost-per-
    hire figure must call this instead of independently re-deriving it, so
    they can never drift apart (copy:both#2: Slide 2 showed "$5,263.16/hire"
    while Slide 6 independently read a raw ``total_projected.cost_per_hire``
    field and showed a different, unexplained "$5,250").

    Returns ``(cph, total_hires)``.
    """
    if not isinstance(budget_alloc, dict):
        return 0.0, 0
    channel_allocs = budget_alloc.get("channel_allocations", {})
    if not isinstance(channel_allocs, dict):
        channel_allocs = {}
    total_hires = sum(
        int(ch.get("projected_hires") or 0)
        for ch in channel_allocs.values()
        if isinstance(ch, dict)
    )
    if total_hires == 0:
        total_proj = budget_alloc.get("total_projected", {})
        if isinstance(total_proj, dict):
            total_hires = int(total_proj.get("hires") or 0)
    metadata = budget_alloc.get("metadata", {})
    total_budget = (
        metadata.get("total_budget") or 0 if isinstance(metadata, dict) else 0
    )
    cph = round(total_budget / total_hires, 2) if total_hires > 0 else 0.0
    return cph, total_hires


def _build_slide_executive_summary(prs: Presentation, data: Dict):
    """Build the Executive Summary slide with hero stat pattern and SCR framework."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    client = data.get("client_name", "Client")
    industry = data.get("industry", "general_entry_level")
    industry_label = data.get("industry_label", industry.replace("_", " ").title())
    locations = data.get("locations") or []
    roles = data.get("roles") or []
    budget = data.get("budget", "TBD")
    work_env = data.get("work_environment", "hybrid")
    goals = _goal_labels(data)
    channels = _selected_channels(data)
    today = datetime.date.today().strftime("%B %d, %Y")

    # Pull synthesized + budget allocation data (from pipeline)
    synthesized = data.get("_synthesized", {})
    if not isinstance(synthesized, dict):
        synthesized = {}
    budget_alloc = data.get("_budget_allocation", {})
    if not isinstance(budget_alloc, dict):
        budget_alloc = {}

    # Extract synthesized sub-sections with safe access
    salary_intel = synthesized.get("salary_intelligence", {})
    if not isinstance(salary_intel, dict):
        salary_intel = {}
    job_market = synthesized.get("job_market_demand", {})
    if not isinstance(job_market, dict):
        job_market = {}

    # Budget allocation sub-sections
    ba_total_projected = budget_alloc.get("total_projected", {})
    if not isinstance(ba_total_projected, dict):
        ba_total_projected = {}
    ba_metadata = budget_alloc.get("metadata", {})
    if not isinstance(ba_metadata, dict):
        ba_metadata = {}

    # S48 FIX: Compute hires from per-channel sum (source of truth) to ensure
    # PPT hero stats match Excel Executive Summary and ROI Projections.
    # copy:both#2: routed through _compute_blended_cph so every slide that
    # shows a blended cost-per-hire derives it identically -- no more
    # independent per-slide re-derivations that can silently disagree.
    _ppt_cph, _ppt_hires_sum = _compute_blended_cph(budget_alloc)

    # Off-white background
    _add_filled_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, OFF_WHITE)

    # Top band
    _add_top_band(slide, "Executive Summary", client.upper())

    # Action title
    role_summary = _cap_roles_for_headline(roles, cap=2)
    loc_count = len(locations)
    loc_text = (
        f"{loc_count} location{'s' if loc_count != 1 else ''}"
        if loc_count > 0
        else "multiple locations"
    )

    # Enhance action text with market temperature if available
    market_temp_str = ""
    try:
        for _role_key, _role_demand in job_market.items():
            if isinstance(_role_demand, dict):
                _temp = _role_demand.get("market_temperature") or ""
                if _temp:
                    market_temp_str = _temp
                    break
    except (AttributeError, TypeError):
        pass

    temp_qualifier = ""
    if market_temp_str:
        temp_qualifier = f" in a {market_temp_str} talent market"

    action_text = (
        f"Nova AI Suite's programmatic strategy targets {role_summary} across "
        f"{loc_text} to optimize "
        f"{client}'s recruitment spend in {industry_label}{temp_qualifier}"
    )
    # O1: measure-then-place. The box is 0.6in tall = 2 lines at 15pt; long
    # client/industry strings (e.g. "Pratt & Whitney New Zealand ... Aerospace &
    # Defense in a hot talent market") wrap to 3 lines and clip behind the
    # SITUATION card at top=1.73in. Autoshrink the headline so it always fits in
    # 2 lines rather than overflowing. Usable width subtracts the text frame's
    # default 0.1in L/R insets.
    _headline_w = 9.9
    _headline_pt = _fit_font_to_lines(
        action_text, _headline_w - 0.2, start_pt=15, max_lines=2
    )
    _add_textbox(
        slide,
        Inches(0.55),
        Inches(0.92),
        Inches(_headline_w),  # narrowed to leave room for the QC chip top-right
        Inches(0.6),
        text=action_text,
        font_size=_headline_pt,
        bold=True,
        color=NAVY,
    )

    # ---- THREE-COLUMN SCR BODY ----
    col_top = Inches(1.65)
    col_height = Inches(3.5)
    col_gap = Inches(0.25)
    accent_bar_w = Inches(0.06)

    col_w = Inches(3.95)
    col1_left = Inches(0.55)
    col2_left = col1_left + col_w + col_gap
    col3_left = col2_left + col_w + col_gap

    # ---- SITUATION (left) ----
    # Light card background
    _add_rounded_rect(slide, col1_left, col_top, col_w, col_height, WHITE)
    _add_filled_rect(slide, col1_left, col_top, accent_bar_w, col_height, BLUE)

    sit_left = col1_left + Inches(0.2)
    sit_w = col_w - Inches(0.25)

    _add_textbox(
        slide,
        sit_left,
        col_top + Inches(0.08),
        sit_w,
        Inches(0.35),
        text="SITUATION",
        font_size=11,
        bold=True,
        color=BLUE,
    )

    body_top = col_top + Inches(0.45)
    _stated_work_label = WORK_ENV_LABELS.get(work_env, work_env.replace("_", " ").title())
    _effective_work_model, _work_model_note = _gold_standard.effective_work_model(
        _stated_work_label, roles
    )
    work_label = _effective_work_model
    role_display = ", ".join(roles[:5]) if roles else "Multiple roles"
    if len(roles) > 5:
        role_display += f" (+{len(roles) - 5} more)"

    # Use total budget from budget engine metadata if available
    total_budget_val = ba_metadata.get("total_budget") or 0
    budget_display_sit = budget
    if total_budget_val and total_budget_val > 0:
        budget_display_sit = _fmt_currency(total_budget_val, compact=True)

    sit_items = [
        ("Industry", industry_label),
        (
            "Locations",
            (
                f"{loc_count} market{'s' if loc_count != 1 else ''}"
                if loc_count > 0
                else "Multiple markets"
            ),
        ),
        ("Target Roles", role_display),
        ("Work Model", work_label),
        ("Budget", budget_display_sit),
    ]
    if _work_model_note:
        sit_items.append(("Work Model Note", _work_model_note))

    # Add market temperature from job_market_demand
    if market_temp_str:
        temp_colors = {
            "hot": "High demand",
            "warm": "Moderate demand",
            "cool": "Balanced",
            "cold": "Low demand",
        }
        sit_items.append(
            (
                "Market Temp.",
                f"{market_temp_str.title()} ({temp_colors.get(market_temp_str, 'N/A')})",
            )
        )

    # v3: Add collar type classification from collar_intelligence
    if _HAS_COLLAR_INTEL and roles:
        try:
            collar_counts = {}
            for role_name in (roles[:8] if isinstance(roles, list) else []):
                r_str = role_name if isinstance(role_name, str) else str(role_name)
                cr = _collar_intel.classify_collar(role=r_str, industry=industry)
                # S4: only count classifications with real confidence -- the
                # classifier's ultimate "no_match" fallback (0.25) is a guess,
                # not a signal about the role, and should never drive a
                # confident-sounding "X dominant (100%)" claim on the deck.
                if cr.get("method") == "no_match" or cr.get("confidence", 0) < 0.4:
                    continue
                ct = cr.get("collar_type", "white_collar")
                collar_counts[ct] = collar_counts.get(ct, 0) + 1
            if collar_counts:
                dominant = max(collar_counts, key=collar_counts.get)
                total_c = sum(collar_counts.values())
                dom_pct = collar_counts[dominant] / total_c * 100
                collar_label = dominant.replace("_", " ").title()
                if dom_pct > 75:
                    sit_items.append(
                        ("Talent Profile", f"{collar_label} dominant ({dom_pct:.0f}%)")
                    )
                elif len(collar_counts) > 1:
                    mix = " / ".join(k.replace("_", " ").title() for k in collar_counts)
                    sit_items.append(("Talent Profile", f"Mixed ({mix})"))
        except Exception:
            pass

    # Add apply rate insight with appropriate framing
    benchmarks = _get_benchmarks(industry, data)

    # Salary Range is added BEFORE CPC Trend / Apply Rate (moved ahead of them,
    # 2026-07-03 finding "SITUATION card's Salary Range row overflows/vanishes
    # on content-heavy plans"): when _autofit_textframe below has to trim
    # trailing paragraphs to make the card fit, it trims from the END of
    # sit_items, so whatever is appended LAST is dropped FIRST. Salary context
    # is at least as decision-relevant as CPC trend/apply-rate framing, so it
    # should not be the row that silently disappears on a content-heavy plan.
    enriched = data.get("_enriched", {})
    salary_data = enriched.get("salary_data", {}) if enriched else {}
    _salary_added = False
    if salary_intel:
        try:
            for _si_role, _si_data in salary_intel.items():
                if isinstance(_si_data, dict):
                    _si_median = _si_data.get("median") or 0
                    _si_min = _si_data.get("min") or 0
                    _si_max = _si_data.get("max") or 0
                    if _si_median and _si_median > 0:
                        salary_str = _format_salary(_si_median)
                        range_str = ""
                        if _si_min > 0 and _si_max > 0:
                            range_str = f" ({_format_salary(_si_min)}-{_format_salary(_si_max)})"
                        sit_items.append(
                            (
                                "Salary Range",
                                f"{salary_str} median{range_str} - {_si_role}",
                            )
                        )
                        _salary_added = True
                        break
        except (AttributeError, TypeError):
            pass
    if not _salary_added and salary_data:
        try:
            first_role = list(salary_data.keys())[0]
            median = salary_data[first_role].get("median") or 0
            if median > 0:
                salary_str = _format_salary(median)
                sit_items.append(
                    ("Salary Benchmark", f"{salary_str} median ({first_role})")
                )
        except (IndexError, KeyError, TypeError):
            pass

    # v3: Add CPC trend direction if available from trend engine
    cpc_trend_str = benchmarks.get("cpc_trend") or ""
    cpc_trend_dir = benchmarks.get("cpc_trend_direction") or ""
    if cpc_trend_str:
        trend_label = {
            "rising": "Rising",
            "falling": "Declining",
            "stable": "Stable",
        }.get(cpc_trend_dir) or ""
        sit_items.append(("CPC Trend", f"{trend_label} {cpc_trend_str}"))

    apply_rate_str = benchmarks.get("apply_rate") or ""
    if apply_rate_str:
        _ar_suffix = _classify_apply_rate_suffix(apply_rate_str)
        if _ar_suffix:
            sit_items.append(("Apply Rate", f"{apply_rate_str} {_ar_suffix}"))

    box2, tf2 = _add_textbox(slide, sit_left, body_top, sit_w, col_height - Inches(0.5))
    tf2.paragraphs[0].space_before = Pt(0)
    tf2.paragraphs[0].space_after = Pt(0)

    first = True
    for label, value in sit_items:
        if first:
            p = tf2.paragraphs[0]
            first = False
        else:
            p = tf2.add_paragraph()
        p.space_before = Pt(2)
        p.space_after = Pt(6)
        p.alignment = PP_ALIGN.LEFT

        run_label = p.add_run()
        run_label.text = f"{label}:  "
        _set_font(run_label, size=10, bold=True, color=DARK_TEXT)

        run_val = p.add_run()
        run_val.text = str(value)
        _set_font(run_val, size=10, bold=False, color=MUTED_TEXT)

    # O1: clamp SITUATION body to the card so the last row never escapes below
    # the card's rounded bottom edge. Fit box height is card bottom
    # (col_top+col_height) minus body_top, less a 0.1in bottom margin.
    _card_body_h = (col_top + col_height - body_top) / 914400 - 0.1
    _autofit_textframe(tf2, sit_w / 914400, _card_body_h)

    # ---- COMPLICATION (middle) ----
    _add_rounded_rect(slide, col2_left, col_top, col_w, col_height, WHITE)
    _add_filled_rect(slide, col2_left, col_top, accent_bar_w, col_height, TEAL)

    comp_left = col2_left + Inches(0.2)
    comp_w = col_w - Inches(0.25)

    _add_textbox(
        slide,
        comp_left,
        col_top + Inches(0.08),
        comp_w,
        Inches(0.35),
        text="COMPLICATION",
        font_size=11,
        bold=True,
        color=TEAL,
    )

    complications = _get_complications(industry)
    box3, tf3 = _add_textbox(
        slide, comp_left, body_top, comp_w, col_height - Inches(0.5)
    )
    tf3.paragraphs[0].space_before = Pt(0)
    tf3.paragraphs[0].space_after = Pt(0)

    for i, item in enumerate(complications):
        if i == 0:
            p = tf3.paragraphs[0]
        else:
            p = tf3.add_paragraph()
        p.space_before = Pt(2)
        p.space_after = Pt(8)
        p.alignment = PP_ALIGN.LEFT

        run_bullet = p.add_run()
        run_bullet.text = "\u25b8  "
        _set_font(run_bullet, size=10, bold=False, color=TEAL)

        run_text = p.add_run()
        run_text.text = str(item) if item is not None else ""
        _set_font(run_text, size=10, bold=False, color=DARK_TEXT)

    # O1: clamp COMPLICATION body to the card bounds (same 3-card pattern).
    _autofit_textframe(tf3, comp_w / 914400, _card_body_h)

    # ---- RESOLUTION (right) ----
    _add_rounded_rect(slide, col3_left, col_top, col_w, col_height, WHITE)
    _add_filled_rect(slide, col3_left, col_top, accent_bar_w, col_height, GREEN)

    res_left = col3_left + Inches(0.2)
    res_w = col_w - Inches(0.25)

    _add_textbox(
        slide,
        res_left,
        col_top + Inches(0.08),
        res_w,
        Inches(0.35),
        text="RESOLUTION",
        font_size=11,
        bold=True,
        color=GREEN,
    )

    box4, tf4 = _add_textbox(slide, res_left, body_top, res_w, col_height - Inches(0.5))
    tf4.paragraphs[0].space_before = Pt(0)
    tf4.paragraphs[0].space_after = Pt(0)

    # Market Thesis -- the WHY this plan will work
    p0 = tf4.paragraphs[0]
    r0 = p0.add_run()
    r0.text = "MARKET THESIS"
    _set_font(r0, size=10, bold=True, color=GREEN)
    p0.space_after = Pt(4)

    # Build thesis from data -- S48: use per-channel-sum hires
    _thesis_parts: list[str] = []
    _thesis_has_lead = False
    if ba_total_projected:
        _proj_h = _ppt_hires_sum
        _proj_cph = _ppt_cph
        if _proj_h > 0:
            _thesis_parts.append(f"This plan projects {int(_proj_h)} hires")
            _thesis_has_lead = True
        if _proj_cph > 0:
            # copy:both#2: fmt_money (never cents) instead of _fmt_currency,
            # which rendered "$5,263.16" whenever the rounded CPH wasn't a
            # whole dollar figure.
            _thesis_parts.append(f"at {_fmt.fmt_money(_proj_cph)}/hire")
    if not _thesis_has_lead:
        # Near-zero-budget edge case (2026-07-03 Gedu review): 0 projected
        # hires means the "This plan projects N hires" clause above never
        # fires, and used to leave the sentence with no subject at all --
        # e.g. a dangling "via 8-channel programmatic strategy." or "in a
        # hot market via 8-channel programmatic strategy." Give the thesis a
        # grammatical lead-in regardless of hire-volume data availability.
        _thesis_parts.append("This plan recommends a targeted approach")
    if market_temp_str:
        temp_map = {
            "hot": "high-demand",
            "warm": "active",
            "cool": "balanced",
            "cold": "buyer's",
        }
        _thesis_parts.append(
            f"in a {temp_map.get(market_temp_str, market_temp_str)} market"
        )
    if len(channels) > 0:
        _thesis_parts.append(f"via {len(channels)}-channel programmatic strategy")

    thesis_text = (
        " ".join(_thesis_parts) + "."
        if _thesis_parts
        else (f"Programmatic multi-channel strategy optimized for {industry_label}.")
    )
    _add_paragraph(
        tf4,
        thesis_text,
        font_size=9,
        color=DARK_TEXT,
        space_before=1,
        space_after=6,
    )

    # Strategy resolution
    _add_paragraph(
        tf4,
        "Nova AI Strategy:",
        font_size=9,
        bold=True,
        color=NAVY,
        space_before=2,
        space_after=2,
    )

    # List every channel in the strategy -- a hardcoded [:5] cap here used to
    # silently omit channels on 6+ channel plans (e.g. an 8-channel plan only
    # showed 5), which read as the checklist contradicting the deck's own
    # "N-channel programmatic strategy" headline a few lines above. O1's
    # autofit/trim-trailing-content mechanism on this card already handles
    # overflow (it shrinks text, then trims the lowest-priority trailing
    # paragraphs -- the publisher line / goals, which sit AFTER this loop --
    # before anything above them would be cut), so listing all channels here
    # is safe rather than silently truncating the ones a client is paying for.
    for ch in list(channels.values()):
        p = tf4.add_paragraph()
        p.space_before = Pt(1)
        p.space_after = Pt(3)
        rb = p.add_run()
        rb.text = "\u2713  "
        _set_font(rb, size=9, bold=False, color=GREEN)
        rt = p.add_run()
        rt.text = ch["label"]
        _set_font(rt, size=9, bold=False, color=DARK_TEXT)

    _total_pubs = _joveo_total_active_publishers(data)
    _pubs_claim = (
        f"across {_total_pubs:,}+ publishers"
        if _total_pubs
        else "across Joveo's publisher network"
    )
    _add_paragraph(
        tf4,
        f"\u2713  ML-optimized bidding {_pubs_claim}",
        font_size=8,
        color=DARK_TEXT,
        space_before=1,
        space_after=3,
    )

    if goals:
        for g in goals[:2]:
            p = tf4.add_paragraph()
            p.space_before = Pt(1)
            p.space_after = Pt(2)
            rb = p.add_run()
            rb.text = "\u25cf  "
            _set_font(rb, size=8, color=BLUE)
            rt = p.add_run()
            rt.text = g
            _set_font(rt, size=8, color=DARK_TEXT)

    # Honest goal-gap one-liner (display_format.parse_hire_goal / goal_gap):
    # only renders when the client stated a hiring GOAL and this plan's
    # projection falls short of it. Framed as a scaling roadmap, not a
    # failure -- matches the fuller goal-gap row on the comparison slide.
    # Placed BEFORE the cited-2026 block (added below): _autofit_textframe
    # trims trailing paragraphs from the END first when a card overflows, so
    # this higher-priority honesty signal must not be added after lower-
    # priority context that's fine to lose first.
    _exec_hire_goal = _fmt.parse_hire_goal(data.get("hire_volume"))
    _exec_goal_gap = _fmt.goal_gap(_ppt_hires_sum, _exec_hire_goal, _ppt_cph)
    if _exec_goal_gap:
        _add_paragraph(
            tf4,
            f"Client goal: {_exec_goal_gap['goal']:,} hires — this plan "
            f"projects {_exec_goal_gap['projected']:,} "
            f"({_exec_goal_gap['pct_of_goal']:.0f}% of goal); scaling path: "
            # copy:both#2: compact fmt_money ("~$2.33M") instead of the raw
            # two-decimal dollar amount ("~$2,331,579.88").
            f"~{_fmt.fmt_money(_exec_goal_gap['additional_budget'], compact=True)} additional",
            font_size=7,
            bold=True,
            color=NAVY,
            space_before=4,
            space_after=1,
        )

    # ---- S82: Cited 2026 market data (parity with Google Slides deck) ----
    # The curated 2026 cited metric + TA-leader quote previously rendered ONLY
    # in the joveo_slides_template (Google Slides tier). When the deck falls
    # back to this python-pptx path, those layers used to be absent. Inject one
    # cited metric here to keep the fallback deck grounded in latest data.
    # Fully isolated: any failure logs and leaves the slide intact.
    try:
        from cited_data_block import build_cited_2026_block, has_cited_content

        _cited = build_cited_2026_block(data)
        if has_cited_content(_cited):
            _cited_bits = []
            for _ml in (_cited.get("metric_lines") or [])[:1]:
                _cited_bits.append(_ml)
            if _cited.get("salary_line") and not _cited_bits:
                _cited_bits.append(_cited["salary_line"])
            if _cited_bits:
                _add_paragraph(
                    tf4,
                    "2026 Market Data:",
                    font_size=8,
                    bold=True,
                    color=NAVY,
                    space_before=4,
                    space_after=1,
                )
                for _cb in _cited_bits:
                    _add_paragraph(
                        tf4,
                        f"\u25aa  {_cb}",
                        font_size=7,
                        color=DARK_TEXT,
                        space_before=1,
                        space_after=1,
                    )
    except Exception as _cited_err:  # pragma: no cover — never break the deck
        logger.debug("Cited 2026 block (exec summary) skipped: %s", _cited_err)

    # O1: clamp RESOLUTION body (market thesis + strategy + cited data) to the
    # card bounds so its final rows never spill below the card / onto the KPI band.
    _autofit_textframe(tf4, res_w / 914400, _card_body_h)

    # ---- HERO STAT METRICS BAR ----
    bar_top = Inches(5.35)
    bar_h = Inches(1.15)

    # Main bar background
    _add_filled_rect(slide, Inches(0.55), bar_top, Inches(12.2), bar_h, NAVY)

    # Teal accent line at top of bar
    _add_filled_rect(slide, Inches(0.55), bar_top, Inches(12.2), Inches(0.04), TEAL)

    # Hero stat: budget (if parseable) or channel count
    budget_display = _format_budget_display(budget)
    hero_value = budget_display if budget_display != budget else str(len(channels))
    hero_label = "Campaign Budget" if budget_display != budget else "Channels Selected"

    # Hero stat on the left -- O1: autoshrink so a long currency-prefixed budget
    # (e.g. "NZ$150,000") never wraps to two lines in the single-line hero slot.
    _hero_pt = _fit_font_single_line(str(hero_value), 3.2 - 0.2, start_pt=36, min_pt=20)
    _add_textbox(
        slide,
        Inches(0.85),
        bar_top + Inches(0.12),
        Inches(3.2),
        Inches(0.52),  # was 0.65 — overlapped the label row below
        text=hero_value,
        font_size=_hero_pt,
        bold=True,
        color=TEAL,
        alignment=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    _add_textbox(
        slide,
        Inches(0.85),
        bar_top + Inches(0.72),
        Inches(3.2),
        Inches(0.3),
        text=hero_label,
        font_size=9,
        bold=False,
        color=LIGHT_MUTED,
        alignment=PP_ALIGN.CENTER,
    )

    # Divider
    _add_filled_rect(
        slide, Inches(4.2), bar_top + Inches(0.2), Inches(0.02), Inches(0.75), TEAL
    )

    # Secondary metrics -- inputs first, then the outcome metrics decision-makers
    # care about. The strip is capped to MAX_SECONDARY and the columns are sized
    # to the available width, so it NEVER marches past the slide edge (the old
    # fixed 1.9in step pushed a 6th/7th KPI 0.7-2.6in off-canvas).
    secondary_metrics = [
        m
        for m in [
            (str(len(channels)), "Channels"),
            (str(loc_count), "Locations") if loc_count > 0 else None,
            (str(len(roles)), "Target Roles") if roles else None,
        ]
        if m is not None
    ]

    # Salary benchmark if available (enrichment)
    if salary_data:
        try:
            first_role = list(salary_data.keys())[0]
            median = salary_data[first_role].get("median") or 0
            if median > 0:
                salary_str = _format_salary(median)
                secondary_metrics.append((salary_str, "Median Salary"))
        except (IndexError, KeyError, TypeError):
            pass

    # Outcome metrics from the budget engine (projected hires, avg CPA)
    # S48: use per-channel-sum hires for consistency
    if ba_total_projected:
        projected_hires = _ppt_hires_sum
        avg_cpa_val = ba_total_projected.get("cost_per_application") or 0
        avg_cph_val = _ppt_cph
        if projected_hires and projected_hires > 0:
            secondary_metrics.append((str(int(projected_hires)), "Projected Hires"))
        if avg_cpa_val and avg_cpa_val > 0:
            secondary_metrics.append((_fmt_currency(avg_cpa_val), "Avg CPA"))
        elif avg_cph_val and avg_cph_val > 0:
            secondary_metrics.append(
                (_fmt_currency(avg_cph_val, compact=True), "Cost/Hire")
            )

    # Market temperature only if there is still room
    if market_temp_str and len(secondary_metrics) < 5:
        secondary_metrics.append((market_temp_str.upper(), "Market Temp."))

    # Hard cap so columns stay readable and on-slide
    MAX_SECONDARY = 5
    secondary_metrics = secondary_metrics[:MAX_SECONDARY]

    # Size columns to the available width (after the hero stat + divider at 4.2in)
    metric_start = Inches(4.45)
    avail_w = Inches(12.75) - metric_start  # right inner margin = 12.75in
    n_secondary = len(secondary_metrics) or 1
    metric_w = Emu(int(avail_w / n_secondary))

    # O1: autoshrink the value font keyed to the rendered string length so a long
    # currency-prefixed value (e.g. "NZ$33.95") never wraps to two lines in this
    # single-line band. This matters more after S3's currency fix, since local
    # currency strings (NZ$/£/etc.) are generally longer than a bare "$" figure.
    _metric_w_in = metric_w / 914400
    _kpi_pt = 24
    for _v, _lbl in secondary_metrics:
        _kpi_pt = min(
            _kpi_pt, _fit_font_single_line(str(_v), _metric_w_in - 0.2, start_pt=24, min_pt=14)
        )

    for i, (value, label) in enumerate(secondary_metrics):
        mx = metric_start + Emu(int(i * metric_w))

        _add_textbox(
            slide,
            mx,
            bar_top + Inches(0.12),
            metric_w,
            Inches(0.55),
            text=value,
            font_size=_kpi_pt,
            bold=True,
            color=WHITE,
            alignment=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        _add_textbox(
            slide,
            mx,
            bar_top + Inches(0.72),
            metric_w,
            Inches(0.3),
            text=label,
            font_size=9,
            bold=False,
            color=LIGHT_MUTED,
            alignment=PP_ALIGN.CENTER,
        )

    # Thin dividers between secondary metrics
    for i in range(1, n_secondary):
        div_x = metric_start + Emu(int(i * metric_w))
        _add_filled_rect(
            slide,
            div_x,
            bar_top + Inches(0.25),
            Inches(0.015),
            Inches(0.65),
            RGBColor(0x1A, 0x45, 0x70),
        )

    # ── Creative Quality Score badge (P1-16) ──
    # S1 (2026-07-03): internal QA artifact -- gated OFF by default so it never
    # reaches a client-facing bundle. Only renders when internal_qc_mode(data)
    # is explicitly enabled (NOVA_INTERNAL_QC env var or data["_internal_preview"]).
    cqs = data.get("_creative_quality_score")
    if (
        _internal_qc_mode(data)
        and cqs
        and isinstance(cqs, dict)
        and cqs.get("score") is not None
        and not cqs.get("degenerate")
    ):
        _cqs_score = cqs.get("score", 0)
        _cqs_grade = cqs.get("grade", "—")
        _cqs_color = (
            GREEN
            if _cqs_score >= 70
            else BLUE if _cqs_score >= 50 else RGBColor(0xCC, 0x33, 0x33)
        )
        # Placed in the top-right CONTENT corner (below the header band, right of
        # the narrowed title, above the SCR cards) so it never collides with the
        # client name in the header band -- the old y=0.18 position overlapped it.
        _add_rounded_rect(
            slide,
            Inches(10.6),
            Inches(0.86),
            Inches(2.15),
            Inches(0.6),
            _cqs_color,
        )
        _add_textbox(
            slide,
            Inches(10.6),
            Inches(0.92),
            Inches(2.15),
            Inches(0.3),
            text=f"Creative QC: {_cqs_score}/100",
            font_size=11,
            bold=True,
            color=WHITE,
            alignment=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        _add_textbox(
            slide,
            Inches(10.6),
            Inches(1.2),
            Inches(2.15),
            Inches(0.22),
            text=f"Grade {_cqs_grade}",
            font_size=9,
            bold=False,
            color=WHITE,
            alignment=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    # Enrichment badge
    _add_enrichment_badge(slide, enriched)

    # Footer
    _add_footer(slide, today)


# ===================================================================
# SLIDE 3 - Section Divider: "Channel Strategy"
# ===================================================================


def _build_slide_divider_channel_strategy(prs: Presentation, data: Dict):
    """Build a full-bleed section divider slide for Channel Strategy."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Full LinkedIn Blue background
    _add_filled_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, BLUE)

    # Teal accent bar at top
    _add_filled_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), TEAL)

    # Teal accent stripe left
    _add_filled_rect(slide, Inches(0.6), Inches(2.8), Inches(2.0), Inches(0.06), TEAL)

    # Section number
    _add_textbox(
        slide,
        Inches(0.6),
        Inches(2.2),
        Inches(3),
        Inches(0.5),
        text="02",
        font_size=18,
        bold=True,
        color=LIGHT_TEAL,
    )

    # Large section title
    _add_textbox(
        slide,
        Inches(0.6),
        Inches(3.1),
        Inches(10),
        Inches(1.5),
        text="Channel Strategy\n& Investment",
        font_size=48,
        bold=True,
        color=WHITE,
    )

    # Subtitle
    _add_textbox(
        slide,
        Inches(0.6),
        Inches(5.0),
        Inches(8),
        Inches(0.5),
        text="Optimized channel mix with programmatic intelligence",
        font_size=16,
        italic=True,
        color=PALE_BLUE,
    )

    # Bottom teal bar
    _add_filled_rect(
        slide, Inches(0), SLIDE_HEIGHT - Inches(0.06), SLIDE_WIDTH, Inches(0.06), TEAL
    )

    # Decorative shapes right side
    _add_oval(
        slide,
        Inches(10.5),
        Inches(1.0),
        Inches(3.5),
        Inches(3.5),
        RGBColor(0x09, 0x58, 0xB0),
    )
    _add_oval(
        slide,
        Inches(11.5),
        Inches(3.5),
        Inches(2.5),
        Inches(2.5),
        RGBColor(0x08, 0x50, 0xA0),
    )


# ===================================================================
# SLIDE 4 - Channel Strategy & Investment with Attribution Diagram
# ===================================================================


def _build_slide_channel_strategy(prs: Presentation, data: Dict):
    """Build the Channel Strategy slide with channel mix bars, benchmarks, and attribution."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    client = data.get("client_name", "Client")
    industry = data.get("industry", "general_entry_level")
    industry_label = data.get("industry_label", industry.replace("_", " ").title())
    channels = _selected_channels(data)
    benchmarks = _get_benchmarks(industry, data)
    today = datetime.date.today().strftime("%B %d, %Y")

    # Pull synthesized + budget allocation data (from pipeline)
    synthesized = data.get("_synthesized", {})
    budget_alloc = data.get("_budget_allocation", {})

    # Off-white background
    _add_filled_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, OFF_WHITE)

    # Top band
    _add_top_band(slide, "CHANNEL STRATEGY & INVESTMENT", today)

    # Action title -- insight-rich with WHY reasoning
    n_cats = len(channels)
    budget_alloc_meta = (
        budget_alloc.get("metadata", {}) if isinstance(budget_alloc, dict) else {}
    )
    total_budget_val = budget_alloc_meta.get("total_budget") or 0
    # S48 FIX: per-channel sum hires for consistency
    _cs_ch_allocs = (
        budget_alloc.get("channel_allocations", {})
        if isinstance(budget_alloc, dict)
        else {}
    )
    if not isinstance(_cs_ch_allocs, dict):
        _cs_ch_allocs = {}
    proj_hires_ch = sum(
        int(ch.get("projected_hires") or 0) for ch in _cs_ch_allocs.values()
    )
    if proj_hires_ch == 0:
        proj_hires_ch = (budget_alloc.get("total_projected", {}) or {}).get(
            "hires"
        ) or 0
    if total_budget_val > 0 and proj_hires_ch > 0:
        action_text = (
            f"{n_cats}-channel strategy allocates {_fmt_currency(total_budget_val, compact=True)} "
            f"to project {int(proj_hires_ch)} hires for {client} in {industry_label} — "
            f"each channel selected for cost-efficiency and audience reach"
        )
    else:
        action_text = (
            f"Optimized {n_cats}-channel mix delivers targeted reach for "
            f"{client}'s {industry_label} hiring priorities with data-driven allocation"
        )
    _add_textbox(
        slide,
        Inches(0.55),
        Inches(0.92),
        Inches(12.2),
        Inches(0.5),
        text=action_text,
        font_size=15,
        bold=True,
        color=NAVY,
    )

    # ==== LEFT: Channel Mix with horizontal bars ====
    left_col_left = Inches(0.55)
    section_top = Inches(1.6)

    # Section header with teal underline
    _add_textbox(
        slide,
        left_col_left,
        section_top,
        Inches(4),
        Inches(0.35),
        text="CHANNEL MIX",
        font_size=11,
        bold=True,
        color=NAVY,
    )
    _add_filled_rect(
        slide,
        left_col_left,
        section_top + Inches(0.33),
        Inches(1.3),
        Inches(0.03),
        TEAL,
    )

    bar_area_top = section_top + Inches(0.5)
    bar_max_w = Inches(3.5)
    bar_h = Inches(0.30)
    bar_spacing = Inches(0.42)
    label_w = Inches(2.3)
    # O1: bar_spacing is finalized below once the channel count is known, so the
    # chart always fits its vertical envelope (above the attribution heading) and
    # the 8th row never collides with / is struck through by that heading.

    # Override channel percentages with real budget allocation if available
    ba_channel_alloc = (
        budget_alloc.get("channel_allocations", {}) if budget_alloc else {}
    )
    if ba_channel_alloc:
        # Map budget engine channel names to display channels
        ba_total_budget = budget_alloc.get("metadata", {}).get("total_budget") or 0
        # S6 fix: reconcile ALL channels' percentages together (largest-remainder
        # rounding) instead of rounding each independently -- independent
        # round(real_pct) per channel is what let an 8-channel mix's displayed
        # percentages drift to 101% (e.g. 31+21+16+13+11+4+3+2 = 101).
        _reconciled_pct = _reconcile_channel_percentages(channels, ba_channel_alloc)
        for ch_key, ch_data in channels.items():
            # Try exact key match, then fuzzy label match
            ba_match = ba_channel_alloc.get(ch_key)
            if not ba_match:
                # Try matching by label (case-insensitive)
                ch_label_lower = (ch_data.get("label") or "").lower()
                for ba_key, ba_val in ba_channel_alloc.items():
                    if isinstance(ba_val, dict):
                        ba_label = ba_val.get("label", ba_key).lower()
                        if (
                            ba_label == ch_label_lower
                            or ba_key.lower() == ch_key.lower()
                        ):
                            ba_match = ba_val
                            break
            if ba_match and isinstance(ba_match, dict):
                real_dollar = ba_match.get("dollar_amount") or 0
                if ch_key in _reconciled_pct:
                    ch_data["pct"] = _reconciled_pct[ch_key]
                if real_dollar > 0:
                    ch_data["_dollar_amount"] = real_dollar

    sorted_channels = sorted(channels.values(), key=lambda c: c["pct"], reverse=True)

    # O1: the CHANNEL CATEGORY ATTRIBUTION heading was pinned to a fixed y=4.85in
    # that silently assumed a ~6-row chart (rows start 2.10in, pitch 0.42in ->
    # 6 rows end at 4.62in). With 8 channels the chart marched to 5.34in and the
    # heading's teal underline struck through the 8th row's label. Fix: hold the
    # heading at its designed position and make the chart occupy a FIXED vertical
    # envelope, tightening the row pitch when the channel count is high so every
    # row (label + bar) lands above the heading with a clear gap. Short charts
    # keep the original 0.42in pitch (min() below), preserving the tested layout.
    _n_bar_rows = len(sorted_channels)
    _attrib_top_in = 4.85  # designed heading position (kept fixed)
    _bar_area_top_in = (section_top + Inches(0.5)) / 914400  # 2.10in
    _bar_h_in = bar_h / 914400  # 0.30in
    _chart_gap_in = 0.20  # clearance between last bar row and the heading
    _chart_envelope_in = _attrib_top_in - _chart_gap_in - _bar_area_top_in
    if _n_bar_rows > 1:
        # pitch that fits N rows (last row needs bar_h, not a full pitch) in the
        # envelope; never larger than the designed 0.42in.
        _fit_pitch_in = (_chart_envelope_in - _bar_h_in) / (_n_bar_rows - 1)
        _pitch_in = min(0.42, _fit_pitch_in)
        bar_spacing = Inches(max(0.30, _pitch_in))  # floor keeps rows readable
    attrib_top = Inches(_attrib_top_in)

    for idx, ch in enumerate(sorted_channels):
        row_y = bar_area_top + idx * bar_spacing

        # Category label (include dollar amount if available from budget engine)
        label_text = ch["label"]
        if ch.get("_dollar_amount"):
            label_text = (
                f"{ch['label']} ({_fmt_currency(ch['_dollar_amount'], compact=True)})"
            )

        _add_textbox(
            slide,
            left_col_left,
            row_y,
            label_w,
            bar_h,
            text=label_text,
            font_size=9,
            bold=True,
            color=DARK_TEXT,
            alignment=PP_ALIGN.RIGHT,
            anchor=MSO_ANCHOR.MIDDLE,
        )

        # Bar
        pct = ch["pct"]
        bar_w_val = bar_max_w * pct / 100
        if bar_w_val < Inches(0.15):
            bar_w_val = Inches(0.15)

        bar_left = left_col_left + label_w + Inches(0.15)
        bar_color = ch.get("color", BLUE)
        _add_rounded_rect(
            slide,
            bar_left,
            row_y + Inches(0.04),
            bar_w_val,
            bar_h - Inches(0.08),
            bar_color,
        )

        # Percentage
        _add_textbox(
            slide,
            bar_left + bar_w_val + Inches(0.08),
            row_y,
            Inches(0.6),
            bar_h,
            text=f"{pct}%",
            font_size=10,
            bold=True,
            color=NAVY,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    # ==== RIGHT TOP: Benchmark Data Table ====
    right_col_left = Inches(7.5)
    right_col_w = Inches(5.3)

    _add_textbox(
        slide,
        right_col_left,
        section_top,
        right_col_w,
        Inches(0.35),
        text="INDUSTRY BENCHMARKS",
        font_size=11,
        bold=True,
        color=NAVY,
    )
    _add_filled_rect(
        slide,
        right_col_left,
        section_top + Inches(0.33),
        Inches(2.0),
        Inches(0.03),
        TEAL,
    )

    table_top = section_top + Inches(0.5)
    table_left = right_col_left
    table_w = Inches(5.1)
    row_h = Inches(0.38)

    # Build benchmark rows with trend annotations (v3)
    cpc_val = benchmarks["cpc"]
    cpa_val = benchmarks["cpa"]
    cpc_trend = benchmarks.get("cpc_trend") or ""
    if cpc_trend:
        cpc_val = f"{cpc_val}  ({cpc_trend})"
    # S3: rows still sourced from a fixed US-calibrated benchmark constant
    # (BENCHMARKS dict / trend_engine / Appcast) get a "US$" marker DIRECTLY ON
    # the figure itself -- never a bare "$" beside this plan's own NZ$/£/etc.
    # figures elsewhere on the slide. Rows overridden by the plan's own live
    # ad-platform / budget-engine data are already in the active currency and
    # need no marker.
    _is_usd_plan = _get_active_currency() == "USD"
    _cpa_label = "Industry CPA"
    _cpc_label = "Industry CPC"
    _cph_label = "Est. Cost-per-Hire"

    # consistency:manpower#2: this row must be the SAME KB "Recruitment
    # Benchmarks" figure the workbook quotes (excel_v2's Executive Summary
    # reads this exact kb['recruitment_benchmarks']['industry_benchmarks']
    # section) -- not this plan's own budget-engine CPH +/- 20%, which
    # Layer 0 of _get_benchmarks() may have substituted into
    # benchmarks["cph"] and which can make the plan look artificially
    # "squarely on-benchmark" against its own derived number. Re-read the
    # KB section directly for this one row.
    _kb_cph_bm = _kb_recruitment_industry_benchmark(industry, data)
    _kb_cph_val = (_kb_cph_bm or {}).get("cph") or ""
    if _kb_cph_val:
        _cph_val = _kb_cph_val
        _cph_is_usd_benchmark = True
    else:
        _cph_val = benchmarks["cph"]
        _cph_is_usd_benchmark = benchmarks.get("cph_is_usd_benchmark", True)

    if not _is_usd_plan:
        if benchmarks.get("cpa_is_usd_benchmark", True):
            cpa_val = _mark_usd(cpa_val)
        if benchmarks.get("cpc_is_usd_benchmark", True):
            cpc_val = _mark_usd(cpc_val)
        if _cph_is_usd_benchmark:
            _cph_val = _mark_usd(_cph_val)
        else:
            # This is the plan's OWN projected cost-per-hire (from the budget
            # engine or salary intelligence), not an external industry
            # constant like the CPA/CPC rows above it -- correctly localized,
            # but relabel it so it doesn't read as an "Industry Benchmark"
            # sourced the same way as its neighbors.
            _cph_label = "Est. Cost-per-Hire (this plan)"
    bench_rows = [
        (_cpa_label, cpa_val),
        (_cpc_label, cpc_val),
        (_cph_label, _cph_val),
        ("Apply Rate", benchmarks["apply_rate"]),
    ]

    # Add real job market data -- prefer synthesized over raw enrichment
    job_market = (
        synthesized.get("job_market_demand", {})
        if isinstance(synthesized, dict) and synthesized
        else {}
    )
    if not job_market:
        enriched = data.get("_enriched", {})
        job_market = enriched.get("job_market", {}) if enriched else {}
    if job_market:
        try:
            for role_name, jm_data in list(job_market.items())[:2]:
                if not isinstance(jm_data, dict) or role_name.startswith("_"):
                    continue
                # Handle both synthesized (total_postings) and raw enriched (posting_count) keys
                posting_count = jm_data.get(
                    "total_postings", jm_data.get("posting_count") or 0
                )
                # S4: fuse_job_market_demand() (data_synthesizer.py) fabricates
                # a generic "job_postings": 75000 fallback when Adzuna/Jooble/
                # search/talent-pool all returned nothing, tagging the source
                # as "Industry Benchmark". Never print that as "Live Postings"
                # -- it is not measured data, and the companion workbook's own
                # Market Demand table for the same role can legitimately show
                # 0/no data for this exact figure (no fabricated stats over
                # empty data).
                _posting_sources = jm_data.get("posting_sources") or []
                _is_fabricated_posting = "Industry Benchmark" in _posting_sources
                avg_sal = jm_data.get("avg_salary") or 0
                if posting_count and posting_count > 0 and not _is_fabricated_posting:
                    bench_rows.append(
                        (
                            f"Live Postings: {role_name}",
                            f"{posting_count:,} active jobs",
                        )
                    )
                if avg_sal and avg_sal > 0:
                    bench_rows.append(
                        (f"Avg Salary: {role_name}", _format_salary(avg_sal))
                    )
                # Synthesized data may have market_temperature
                market_temp = jm_data.get("market_temperature") or ""
                if market_temp and isinstance(market_temp, str):
                    bench_rows.append(
                        (f"Market Temp: {role_name}", market_temp.title())
                    )
        except (TypeError, AttributeError):
            pass

    # Add real ad platform analysis data from synthesized pipeline
    ad_plat = synthesized.get("ad_platform_analysis", {}) if synthesized else {}
    if ad_plat:
        try:
            for plat_name, plat_data in list(ad_plat.items())[:5]:
                if not isinstance(plat_data, dict) or plat_name.startswith("_"):
                    continue
                plat_label = plat_data.get("platform_name", _smart_title(plat_name))
                plat_cpc = plat_data.get("CPC", plat_data.get("cpc") or 0)
                plat_cpa = plat_data.get("CPA", plat_data.get("cpa") or 0)
                plat_reach = plat_data.get("estimated_reach") or 0
                fit_score = plat_data.get("fit_score") or 0
                if plat_cpc and plat_cpc > 0:
                    bench_rows.append((f"{plat_label} CPC", _fmt_currency(plat_cpc)))
                if plat_cpa and plat_cpa > 0:
                    bench_rows.append((f"{plat_label} CPA", _fmt_currency(plat_cpa)))
                if plat_reach and plat_reach > 0:
                    bench_rows.append((f"{plat_label} Est. Reach", f"{plat_reach:,}"))
                if fit_score and fit_score > 0:
                    bench_rows.append(
                        (f"{plat_label} Fit Score", _fmt_pct(fit_score, decimals=0))
                    )
                # Deep intelligence data (91-platform KB enrichment)
                deep = plat_data.get("deep_intelligence", {})
                if isinstance(deep, dict) and deep:
                    visitors = deep.get("monthly_visitors")
                    if visitors:
                        bench_rows.append(
                            (f"{plat_label} Monthly Visitors", str(visitors))
                        )
                    best_for = deep.get("best_for") or []
                    if isinstance(best_for, list) and best_for:
                        bench_rows.append(
                            (
                                f"{plat_label} Best For",
                                ", ".join(str(b) for b in best_for[:3]),
                            )
                        )
        except (TypeError, AttributeError):
            pass

        # Programmatic insights from supply ecosystem KB
        prog_insights = ad_plat.get("_programmatic_insights", {})
        if isinstance(prog_insights, dict) and prog_insights:
            try:
                bidding = prog_insights.get("bidding_models", {})
                if isinstance(bidding, dict) and bidding:
                    for bk, bv in list(bidding.items())[:2]:
                        label = str(bk).replace("_", " ").title()
                        if isinstance(bv, dict):
                            desc = bv.get(
                                "description", str(next(iter(bv.values()), ""))
                            )
                        else:
                            desc = str(bv)
                        bench_rows.append((f"Bidding: {label}", str(desc)[:40]))
            except (TypeError, AttributeError):
                pass

    # ── LinkedIn Intelligence (SlotOps 108K dataset) ──
    li_intel = (data.get("_gold_standard") or {}).get("linkedin_intelligence", {})
    if not li_intel:
        li_intel = data.get("_slotops_linkedin_benchmarks", {})
    if li_intel:
        country_ar = li_intel.get("country_apply_rate", {})
        ea_ats_data = li_intel.get("ea_vs_ats", {})
        li_country = li_intel.get("country", "")
        li_sample = li_intel.get("sample_size", 0)
        if country_ar.get("avg"):
            bench_rows.append(
                (
                    f"LinkedIn Apply Rate ({li_country})",
                    f"{country_ar['avg']:.1f}% avg ({li_sample:,} jobs)",
                )
            )
        if ea_ats_data.get("easy_apply_rate") and ea_ats_data.get("ats_rate"):
            lift = ea_ats_data.get("lift_factor", 0)
            bench_rows.append(
                (
                    "LinkedIn Easy Apply vs ATS",
                    f"{ea_ats_data['easy_apply_rate']:.1f}% vs "
                    f"{ea_ats_data['ats_rate']:.1f}% ({lift:.1f}x lift)",
                )
            )
        best_days = li_intel.get("best_posting_days", [])
        if best_days:
            bench_rows.append(("LinkedIn Best Posting Days", ", ".join(best_days)))

    # Has ad platform data - use for 3-column table header
    has_ad_plat_data = bool(ad_plat)

    # Cap rows so the table fits ABOVE the attribution band (attrib_top=4.85in)
    # and never runs off the bottom / overlaps the category cards. The data-driven
    # row list could reach 13+ rows (ad-platform + LinkedIn intel); the reference
    # deck shows ~5 benchmark rows. The fuller breakdown lives in the workbook.
    # max rows = floor((attrib_top - clearance - source - table_top)/row_h) - header
    # O1: use the dynamically-computed attribution top (row-count aware) so the
    # right-hand benchmark table is capped to stay above wherever the heading
    # actually lands, not a stale fixed 4.85in.
    _max_bench_rows = max(
        4,
        int((_attrib_top_in - 0.35 - (table_top + row_h) / 914400) / (row_h / 914400)),
    )
    if len(bench_rows) > _max_bench_rows:
        bench_rows = bench_rows[:_max_bench_rows]

    # Table header
    _add_filled_rect(slide, table_left, table_top, table_w, row_h, NAVY)
    _add_textbox(
        slide,
        table_left + Inches(0.15),
        table_top,
        Inches(2.2),
        row_h,
        text="Metric",
        font_size=9,
        bold=True,
        color=WHITE,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    _add_textbox(
        slide,
        table_left + Inches(2.4),
        table_top,
        Inches(2.5),
        row_h,
        text=f"{industry_label} Range",
        font_size=9,
        bold=True,
        color=WHITE,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # v3: Determine value color based on confidence level
    bench_conf = benchmarks.get("confidence", "curated")
    conf_value_color, _conf_label = _confidence_color(bench_conf)

    for i, (metric, value) in enumerate(bench_rows):
        ry = table_top + row_h * (i + 1)
        bg = WHITE if i % 2 == 0 else RGBColor(0xF8, 0xF6, 0xF3)
        _add_filled_rect(slide, table_left, ry, table_w, row_h, bg)
        # Thin ruled line between rows (McKinsey/Bain style)
        _add_filled_rect(
            slide,
            table_left,
            ry + row_h - Inches(0.01),
            table_w,
            Inches(0.01),
            WARM_GRAY,
        )

        _add_textbox(
            slide,
            table_left + Inches(0.15),
            ry,
            Inches(2.2),
            row_h,
            text=metric,
            font_size=9,
            bold=True,
            color=DARK_TEXT,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        # Use confidence-colored text for industry benchmark rows (first 4)
        # Live API rows use green, trend engine uses blue, curated uses navy
        val_color = conf_value_color if i < 4 else NAVY
        _add_textbox(
            slide,
            table_left + Inches(2.4),
            ry,
            Inches(2.5),
            row_h,
            text=str(value),
            font_size=10,
            bold=True,
            color=val_color,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    # Source - adjust position based on actual number of rows
    source_top = table_top + row_h * (len(bench_rows) + 1) + Inches(0.05)
    source_text = f"Sources: Industry benchmarks {datetime.date.today().year}, validated recruitment data"
    if ad_plat:
        source_text += ", Nova AI Suite Ad Platform Intelligence"
    _add_textbox(
        slide,
        table_left,
        source_top,
        table_w,
        Inches(0.2),
        text=source_text,
        font_size=7,
        italic=True,
        color=MUTED_TEXT,
    )

    # ==== CHANNEL ATTRIBUTION DIAGRAM (bottom area) ====
    # attrib_top computed above from the actual channel-mix row count (O1).
    _add_textbox(
        slide,
        Inches(0.55),
        attrib_top,
        Inches(12.2),
        Inches(0.35),
        text="CHANNEL CATEGORY ATTRIBUTION",
        font_size=11,
        bold=True,
        color=NAVY,
    )
    _add_filled_rect(
        slide, Inches(0.55), attrib_top + Inches(0.33), Inches(2.8), Inches(0.03), TEAL
    )

    # Build category groups
    cat_groups = _channel_categories_grouped(channels)

    # Attribution category boxes
    cat_colors = {
        "Programmatic": (NAVY, WHITE),
        "Job Boards": (BLUE, WHITE),
        "Social": (SKY_BLUE, NAVY),
        "Employer Brand": (TEAL, NAVY),
        "Other": (MEDIUM_BLUE, WHITE),
    }

    box_top = attrib_top + Inches(0.5)
    box_h = Inches(1.2)
    total_available_w = Inches(12.2)
    n_groups = len(cat_groups)

    if n_groups > 0:
        box_gap = Inches(0.15)
        box_w = (
            (total_available_w - box_gap * (n_groups - 1)) / n_groups
            if n_groups > 1
            else total_available_w
        )
        overlap_w = Inches(0.3)  # visual overlap zone

        for gi, (cat_name, cat_channels) in enumerate(cat_groups.items()):
            bx = Inches(0.55) + gi * (box_w + box_gap)
            bg_color, text_color = cat_colors.get(cat_name, (MEDIUM_BLUE, WHITE))

            # Category card
            _add_rounded_rect(slide, bx, box_top, box_w, box_h, bg_color)

            # Category name
            _add_textbox(
                slide,
                bx + Inches(0.15),
                box_top + Inches(0.08),
                box_w - Inches(0.3),
                Inches(0.3),
                text=cat_name.upper(),
                font_size=10,
                bold=True,
                color=text_color,
            )

            # Total percentage for this category
            cat_pct = sum(c["pct"] for c in cat_channels)
            _add_textbox(
                slide,
                bx + Inches(0.15),
                box_top + Inches(0.35),
                box_w - Inches(0.3),
                Inches(0.35),
                text=f"{cat_pct}%",
                font_size=22,
                bold=True,
                color=text_color,
            )

            # Channel list
            ch_list = ", ".join(c["label"] for c in cat_channels)
            _add_textbox(
                slide,
                bx + Inches(0.15),
                box_top + Inches(0.72),
                box_w - Inches(0.3),
                Inches(0.42),
                text=ch_list,
                font_size=7,
                color=text_color,
            )

        # Overlap connectors between categories (teal diamonds)
        for gi in range(n_groups - 1):
            connector_x = (
                Inches(0.55) + (gi + 1) * (box_w + box_gap) - box_gap / 2 - Inches(0.12)
            )
            connector_y = box_top + box_h / 2 - Inches(0.12)
            diamond = slide.shapes.add_shape(
                MSO_SHAPE.DIAMOND, connector_x, connector_y, Inches(0.24), Inches(0.24)
            )
            diamond.fill.solid()
            diamond.fill.fore_color.rgb = TEAL
            diamond.line.fill.background()

    # v3: Data sources footnote with confidence indicator
    _add_data_sources_footnote(slide, data, benchmarks)

    # Footer
    _add_footer(slide, today)


# ===================================================================
# SLIDE 5 - Quality & ROI Outcomes Grid
# ===================================================================


def _build_slide_quality_outcomes(prs: Presentation, data: Dict):
    """Build the Quality Outcomes grid slide with 4-quadrant metrics."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    client = data.get("client_name", "Client")
    industry = data.get("industry", "general_entry_level")
    channels = _selected_channels(data)
    budget = data.get("budget", "TBD")
    roles = data.get("roles") or []
    locations = data.get("locations") or []
    today = datetime.date.today().strftime("%B %d, %Y")

    # Pull synthesized + budget allocation data (from pipeline)
    synthesized = data.get("_synthesized", {})
    budget_alloc = data.get("_budget_allocation", {})

    # Off-white background
    _add_filled_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, OFF_WHITE)

    # Top band
    _add_top_band(slide, "QUALITY & ROI PROJECTIONS", today)

    # Action title
    n_channels = len(channels)
    action_text = (
        f"Projected quality outcomes across {n_channels} optimized channels "
        f"for {client}'s programmatic media plan"
    )
    _add_textbox(
        slide,
        Inches(0.55),
        Inches(0.92),
        Inches(12.2),
        Inches(0.5),
        text=action_text,
        font_size=15,
        bold=True,
        color=NAVY,
    )

    # NOTE: a centered hero "budget" card used to sit at y=1.55-2.85, directly on
    # top of the 5-card projections row below (y=1.9-2.9). It was redundant (the
    # budget already appears on the cover, exec summary, and budget slides), so it
    # was removed -- the 5-card projections row is now the clean top section.

    # ---- CAMPAIGN PROJECTIONS SUMMARY (5-card row) ----
    ba_total_proj = budget_alloc.get("total_projected", {}) if budget_alloc else {}
    if not isinstance(ba_total_proj, dict):
        ba_total_proj = {}
    ba_metadata_qo = budget_alloc.get("metadata", {}) if budget_alloc else {}
    if not isinstance(ba_metadata_qo, dict):
        ba_metadata_qo = {}

    proj_clicks = ba_total_proj.get("clicks") or 0
    proj_apps = ba_total_proj.get("applications") or 0
    # S48 FIX: Compute hires from per-channel sum for consistency
    _qo_ch_allocs = budget_alloc.get("channel_allocations", {}) if budget_alloc else {}
    if not isinstance(_qo_ch_allocs, dict):
        _qo_ch_allocs = {}
    projected_hires = sum(
        int(ch.get("projected_hires") or 0) for ch in _qo_ch_allocs.values()
    )
    if projected_hires == 0:
        projected_hires = ba_total_proj.get("hires") or 0
    real_avg_cpa = ba_total_proj.get("cost_per_application") or 0
    _qo_budget = ba_metadata_qo.get("total_budget") or 0
    ba_avg_cph = (
        round(_qo_budget / max(projected_hires, 1), 2)
        if projected_hires > 0
        else (ba_total_proj.get("cost_per_hire") or 0)
    )

    benchmarks = _get_benchmarks(industry, data)
    cpa_str = benchmarks.get("cpa", "$25")
    try:
        cpa_nums = re.findall(r"[\d.]+", cpa_str.replace(",", ""))
        benchmark_avg_cpa = (
            sum(float(x) for x in cpa_nums) / len(cpa_nums) if cpa_nums else 25
        )
    except Exception:
        benchmark_avg_cpa = 25
    avg_cpa = real_avg_cpa if real_avg_cpa and real_avg_cpa > 0 else benchmark_avg_cpa
    efficiency_improvement = min(35, max(15, round(100 / avg_cpa * 5)))

    enriched = data.get("_enriched", {})
    salary_data = enriched.get("salary_data", {}) if enriched else {}

    # Section label
    _add_textbox(
        slide,
        Inches(0.55),
        Inches(1.5),
        Inches(5),
        Inches(0.28),
        text="CAMPAIGN PROJECTIONS SUMMARY",
        font_size=10,
        bold=True,
        color=BLUE,
    )
    _add_filled_rect(slide, Inches(0.55), Inches(1.76), Inches(2.5), Inches(0.03), TEAL)

    # 5-metric summary cards
    summary_top = Inches(1.9)
    summary_h = Inches(1.0)
    card_w = Inches(2.3)
    card_gap = Inches(0.12)
    card_start_x = Inches(0.55)

    summary_metrics = [
        {
            "value": f"{proj_clicks:,}" if proj_clicks > 0 else "--",
            "label": "Projected Clicks",
            "accent": BLUE,
        },
        {
            "value": f"{int(proj_apps):,}" if proj_apps > 0 else "--",
            "label": "Projected Applications",
            "accent": TEAL,
        },
        {
            "value": f"{int(projected_hires):,}" if projected_hires > 0 else "--",
            "label": "Projected Hires",
            "accent": GREEN,
        },
        {
            "value": _fmt_currency(avg_cpa) if avg_cpa > 0 else "--",
            "label": "Avg CPA",
            "accent": RGBColor(0xED, 0x7D, 0x31),
        },
        {
            "value": _fmt_currency(ba_avg_cph) if ba_avg_cph > 0 else "--",
            "label": "Avg Cost/Hire",
            "accent": NAVY,
        },
    ]

    for si, sm in enumerate(summary_metrics):
        sx = card_start_x + si * (card_w + card_gap)
        _add_rounded_rect(slide, sx, summary_top, card_w, summary_h, WHITE)
        _add_filled_rect(slide, sx, summary_top, card_w, Inches(0.04), sm["accent"])
        _add_textbox(
            slide,
            sx,
            summary_top + Inches(0.1),
            card_w,
            Inches(0.5),
            text=sm["value"],
            font_size=26,
            bold=True,
            color=sm["accent"],
            alignment=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        _add_textbox(
            slide,
            sx,
            summary_top + Inches(0.65),
            card_w,
            Inches(0.25),
            text=sm["label"],
            font_size=9,
            bold=False,
            color=MUTED_TEXT,
            alignment=PP_ALIGN.CENTER,
        )

    # ---- CHANNEL BREAKDOWN TABLE ----
    ch_table_top = Inches(3.15)
    ch_table_left = Inches(0.55)
    ch_table_w = Inches(12.2)

    # Section label
    _add_textbox(
        slide,
        ch_table_left,
        ch_table_top,
        Inches(5),
        Inches(0.28),
        text="CHANNEL-BY-CHANNEL PROJECTIONS",
        font_size=10,
        bold=True,
        color=BLUE,
    )

    # Get channel allocations
    ba_channel_alloc_qo = (
        budget_alloc.get("channel_allocations", {}) if budget_alloc else {}
    )
    if not isinstance(ba_channel_alloc_qo, dict):
        ba_channel_alloc_qo = {}

    # Build display data for channels
    ch_display_list = []
    for ch_key, ch_data in ba_channel_alloc_qo.items():
        if not isinstance(ch_data, dict):
            continue
        ch_display_list.append(
            {
                "label": CHANNEL_ALLOC.get(ch_key, {}).get("label")
                or ch_key.replace("_", " ").title(),
                "budget": ch_data.get("dollar_amount", ch_data.get("dollars") or 0),
                "clicks": ch_data.get("projected_clicks") or 0,
                "apps": ch_data.get("projected_applications") or 0,
                "hires": ch_data.get("projected_hires") or 0,
                "cpa": ch_data.get("cpa", ch_data.get("cost_per_application") or 0),
            }
        )

    # If no budget engine data, try to build from channels dict
    if not ch_display_list:
        ba_total_budget_qo = ba_metadata_qo.get("total_budget") or 0
        for ch_key, ch_data in channels.items():
            ch_pct = ch_data.get("pct") or 0
            ch_dollars = (
                ba_total_budget_qo * ch_pct / 100.0 if ba_total_budget_qo > 0 else 0
            )
            ch_display_list.append(
                {
                    "label": ch_data.get("label", ch_key.replace("_", " ").title()),
                    "budget": ch_dollars,
                    "clicks": 0,
                    "apps": 0,
                    "hires": 0,
                    "cpa": 0,
                }
            )

    # Sort by budget descending, take top 5
    ch_display_list.sort(key=lambda c: c.get("budget") or 0, reverse=True)
    ch_display_top5 = ch_display_list[:5]

    # Table header row
    header_y = ch_table_top + Inches(0.32)
    row_h = Inches(0.34)
    col_widths_qo = [
        Inches(3.0),
        Inches(2.0),
        Inches(1.8),
        Inches(1.8),
        Inches(1.8),
        Inches(1.8),
    ]
    col_headers_qo = ["Channel", "Budget", "Clicks", "Applications", "Hires", "CPA"]
    col_aligns_qo = [
        PP_ALIGN.LEFT,
        PP_ALIGN.CENTER,
        PP_ALIGN.CENTER,
        PP_ALIGN.CENTER,
        PP_ALIGN.CENTER,
        PP_ALIGN.CENTER,
    ]

    # Header background
    _add_filled_rect(slide, ch_table_left, header_y, ch_table_w, row_h, NAVY)
    cx = ch_table_left
    for ci, (header, cw) in enumerate(zip(col_headers_qo, col_widths_qo)):
        _add_textbox(
            slide,
            cx + Inches(0.1),
            header_y,
            cw - Inches(0.1),
            row_h,
            text=header,
            font_size=9,
            bold=True,
            color=WHITE,
            alignment=col_aligns_qo[ci],
            anchor=MSO_ANCHOR.MIDDLE,
        )
        cx += cw

    # Data rows (top 5 channels)
    for ri, ch in enumerate(ch_display_top5):
        row_y = header_y + row_h + ri * row_h
        row_bg = WHITE if ri % 2 == 0 else RGBColor(0xF5, 0xF5, 0xF3)
        _add_filled_rect(slide, ch_table_left, row_y, ch_table_w, row_h, row_bg)

        row_values = [
            ch["label"],
            _fmt_currency(ch["budget"]) if ch["budget"] > 0 else "--",
            f"{int(ch['clicks']):,}" if ch["clicks"] > 0 else "--",
            f"{int(ch['apps']):,}" if ch["apps"] > 0 else "--",
            f"{int(ch['hires']):,}" if ch["hires"] > 0 else "--",
            _fmt_currency(ch["cpa"]) if ch["cpa"] > 0 else "--",
        ]

        cx = ch_table_left
        for ci, (val, cw) in enumerate(zip(row_values, col_widths_qo)):
            left_pad = Inches(0.15) if ci == 0 else Inches(0.1)
            _add_textbox(
                slide,
                cx + left_pad,
                row_y,
                cw - left_pad,
                row_h,
                text=val,
                font_size=9,
                bold=(ci == 0),
                color=DARK_TEXT,
                alignment=col_aligns_qo[ci],
                anchor=MSO_ANCHOR.MIDDLE,
            )
            cx += cw

    # ---- BUDGET REALITY CHECK (if budget is insufficient) ----
    _suff_data = budget_alloc.get("sufficiency", {}) if budget_alloc else {}
    if not isinstance(_suff_data, dict):
        _suff_data = {}
    _budget_reality = (
        budget_alloc.get("budget_reality_check", {}) if budget_alloc else {}
    )
    if not isinstance(_budget_reality, dict):
        _budget_reality = {}

    _is_critical = False
    _reality_message = ""

    # Check budget_reality_check first (if another agent added it)
    if _budget_reality:
        _feas_tier = _budget_reality.get("feasibility_tier") or ""
        if _feas_tier in ("impossible", "severely_underfunded"):
            _is_critical = True
            _reality_message = _budget_reality.get("feasibility_message") or ""
            if not _reality_message:
                _reality_message = (
                    f"Budget is {_budget_reality.get('feasibility_label', 'severely underfunded')}. "
                    f"Budget per hire: {_fmt_currency(_budget_reality.get('budget_per_hire') or 0)} vs. "
                    f"industry avg: {_fmt_currency(_budget_reality.get('industry_avg_cph') or 0)}."
                )
    # Fall back to sufficiency data
    elif _suff_data and not _suff_data.get("sufficient", True):
        _is_critical = True
        _gap = _suff_data.get("gap_amount") or 0
        _avg_cph_suff = _suff_data.get("industry_avg_cost_per_hire") or 0
        _bpo = _suff_data.get("budget_per_opening") or 0
        _reality_message = (
            f"Budget per opening (${_bpo:,.0f}) is below industry average "
            f"cost-per-hire (${_avg_cph_suff:,.0f}). "
        )
        if _gap > 0:
            _reality_message += (
                f"An additional ${_gap:,.0f} is recommended to meet all hiring targets."
            )

    # Position for reality check or insight callout -- dynamically BELOW the
    # actual channel table bottom so it never overlaps the table rows (the old
    # fixed 5.15in sat on top of a 5-row table that ended at ~5.51in).
    _qo_table_bottom = header_y + row_h * (len(ch_display_top5) + 1)
    bottom_section_top = _qo_table_bottom + Inches(0.16)

    if _is_critical and _reality_message:
        # Red callout box for budget reality check
        reality_top = bottom_section_top
        reality_h = Inches(0.7)
        RED_BG = RGBColor(0xFD, 0xE8, 0xE8)
        RED_ACCENT = RGBColor(0xC6, 0x28, 0x28)
        _add_rounded_rect(
            slide, Inches(0.55), reality_top, Inches(12.2), reality_h, RED_BG
        )
        _add_filled_rect(
            slide, Inches(0.55), reality_top, Inches(0.06), reality_h, RED_ACCENT
        )

        # Badge
        _add_rounded_rect(
            slide,
            Inches(0.85),
            reality_top + Inches(0.17),
            Inches(2.0),
            Inches(0.35),
            RED_ACCENT,
        )
        _add_textbox(
            slide,
            Inches(0.85),
            reality_top + Inches(0.17),
            Inches(2.0),
            Inches(0.35),
            text="BUDGET REALITY CHECK",
            font_size=8,
            bold=True,
            color=WHITE,
            alignment=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )

        _add_textbox(
            slide,
            Inches(3.1),
            reality_top + Inches(0.1),
            Inches(9.4),
            reality_h - Inches(0.15),
            text=_reality_message,
            font_size=10,
            bold=False,
            color=RGBColor(0xC6, 0x28, 0x28),
        )

        # Shift insight callout below
        insight_top = reality_top + reality_h + Inches(0.1)
    else:
        insight_top = bottom_section_top

    # ---- KEY INSIGHT CALLOUT BOX ----
    insight_h = Inches(0.85)
    _add_rounded_rect(
        slide, Inches(0.55), insight_top, Inches(12.2), insight_h, PALE_TEAL
    )
    _add_filled_rect(slide, Inches(0.55), insight_top, Inches(0.06), insight_h, TEAL)

    # Insight icon/badge
    _add_rounded_rect(
        slide, Inches(0.85), insight_top + Inches(0.15), Inches(1.0), Inches(0.35), TEAL
    )
    _add_textbox(
        slide,
        Inches(0.85),
        insight_top + Inches(0.15),
        Inches(1.0),
        Inches(0.35),
        text="KEY INSIGHT",
        font_size=8,
        bold=True,
        color=WHITE,
        alignment=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # Build insight text using real data when available
    if real_avg_cpa and real_avg_cpa > 0:
        insight_text = (
            f"Nova AI Suite's programmatic approach distributes {client}'s budget across "
            f"{n_channels} optimized channels with ML-driven bid management, "
            f"projecting ${real_avg_cpa:,.0f} avg CPA (vs. industry benchmark ${benchmark_avg_cpa:.0f}). "
            f"Quality-focused optimization (CPQA) ensures spend is directed toward "
            f"candidates most likely to apply and convert."
        )
    else:
        insight_text = (
            f"Nova AI Suite's programmatic approach distributes {client}'s budget across "
            f"{n_channels} optimized channels with ML-driven bid management, "
            f"projecting {efficiency_improvement}% CPA improvement over manual posting. "
            f"Quality-focused optimization (CPQA) ensures spend is directed toward "
            f"candidates most likely to apply and convert."
        )

    # Append projected hires if available from budget allocation
    if projected_hires and projected_hires > 0:
        total_apps_insight = ba_total_proj.get("applications") or 0
        if total_apps_insight and total_apps_insight > 0:
            insight_text += (
                f" Budget engine projects {int(total_apps_insight):,} applications and "
                f"{int(projected_hires):,} hires from the allocated investment."
            )
        else:
            insight_text += (
                f" Budget engine projects {int(projected_hires):,} hires "
                f"from the allocated investment."
            )

    # Append salary insight if enrichment data is available
    if salary_data:
        try:
            first_role = list(salary_data.keys())[0]
            median = salary_data[first_role].get("median") or 0
            if median > 0:
                insight_text += (
                    f" Market salary data shows {_format_salary(median)} "
                    f"median for {first_role}, enabling precise budget calibration."
                )
        except (IndexError, KeyError, TypeError):
            pass

    _add_textbox(
        slide,
        Inches(2.1),
        insight_top + Inches(0.08),
        Inches(10.4),
        insight_h - Inches(0.15),
        text=_trunc_word(insight_text, 500),
        font_size=9,
        color=DARK_TEXT,
    )

    # Enrichment badge
    _add_enrichment_badge(slide, enriched)

    # Footer
    _add_footer(slide, today)


# ===================================================================
# SLIDE 6 - Budget Allocation & Projections
# ===================================================================


def _build_slide_budget_allocation(prs: Presentation, data: Dict):
    """Build a dedicated Budget Allocation slide showing dollar breakdown per channel,
    projected applications, projected hires, and ROI projections.

    This slide is only added when real budget allocation data is available from the
    budget engine. It provides the financial transparency Fortune 500 clients expect.
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    client = data.get("client_name", "Client")
    industry = data.get("industry", "general_entry_level")
    channels = _selected_channels(data)
    budget = data.get("budget", "TBD")
    today = datetime.date.today().strftime("%B %d, %Y")

    budget_alloc = data.get("_budget_allocation", {})
    if not isinstance(budget_alloc, dict):
        budget_alloc = {}
    ba_total_proj = budget_alloc.get("total_projected", {})
    if not isinstance(ba_total_proj, dict):
        ba_total_proj = {}
    ba_channel_alloc = budget_alloc.get("channel_allocations", {})
    if not isinstance(ba_channel_alloc, dict):
        ba_channel_alloc = {}
    ba_metadata = budget_alloc.get("metadata", {})
    if not isinstance(ba_metadata, dict):
        ba_metadata = {}
    ba_total_budget = ba_metadata.get("total_budget") or 0

    enriched = data.get("_enriched", {})

    # Off-white background
    _add_filled_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, OFF_WHITE)

    # Top band
    _add_top_band(slide, "BUDGET ALLOCATION & PROJECTIONS", today)

    # Action title
    n_channels = len(channels)
    action_text = (
        f"Investment breakdown across {n_channels} channels "
        f"with projected outcomes for {client}"
    )
    _add_textbox(
        slide,
        Inches(0.55),
        Inches(0.92),
        Inches(12.2),
        Inches(0.45),
        text=action_text,
        font_size=15,
        bold=True,
        color=NAVY,
    )

    # ---- HERO STATS ROW (3 cards) ----
    hero_top = Inches(1.5)
    hero_h = Inches(1.1)
    hero_w = Inches(3.8)
    hero_gap = Inches(0.35)
    hero_start_x = Inches(0.55)

    # Total Investment
    _cur = _cur_symbol()
    total_display = (
        f"{_cur}{ba_total_budget:,.0f}"
        if ba_total_budget > 0
        else _format_budget_display(budget)
    )

    # Projected Applications
    proj_apps = ba_total_proj.get("applications") or 0
    apps_display = f"{int(proj_apps):,}" if proj_apps and proj_apps > 0 else "--"

    # Projected Hires -- S48: per-channel sum for consistency
    proj_hires = sum(
        int(ch.get("projected_hires") or 0) for ch in ba_channel_alloc.values()
    )
    if proj_hires == 0:
        proj_hires = ba_total_proj.get("hires") or 0
    hires_display = f"{int(proj_hires):,}" if proj_hires and proj_hires > 0 else "--"

    hero_cards = [
        {"value": total_display, "label": "Total Investment", "accent": BLUE},
        {"value": apps_display, "label": "Projected Applications", "accent": TEAL},
        {"value": hires_display, "label": "Projected Hires", "accent": GREEN},
    ]

    for hi, hc in enumerate(hero_cards):
        hx = hero_start_x + hi * (hero_w + hero_gap)
        _add_rounded_rect(slide, hx, hero_top, hero_w, hero_h, WHITE)
        _add_filled_rect(slide, hx, hero_top, hero_w, Inches(0.05), hc["accent"])
        _add_textbox(
            slide,
            hx,
            hero_top + Inches(0.12),
            hero_w,
            Inches(0.6),
            text=hc["value"],
            font_size=34,
            bold=True,
            color=hc["accent"],
            alignment=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        _add_textbox(
            slide,
            hx,
            hero_top + Inches(0.72),
            hero_w,
            Inches(0.3),
            text=hc["label"],
            font_size=11,
            bold=False,
            color=MUTED_TEXT,
            alignment=PP_ALIGN.CENTER,
        )

    # ---- CHANNEL BREAKDOWN TABLE ----
    table_top = Inches(2.85)
    table_left = Inches(0.55)
    table_w = Inches(12.2)

    # Section label
    _add_textbox(
        slide,
        table_left,
        table_top,
        Inches(5),
        Inches(0.3),
        text="CHANNEL-BY-CHANNEL BREAKDOWN",
        font_size=10,
        bold=True,
        color=BLUE,
    )

    # S6 fix: reconcile ALL channels' percentages together (largest-remainder
    # rounding) instead of rounding each independently -- this is what let an
    # 8-channel plan's Total row print "101%" while the underlying dollars
    # summed to exactly the stated budget.
    _reconciled_pct = _reconcile_channel_percentages(channels, ba_channel_alloc)

    # Map budget engine channel data onto our display channels
    display_channels = []
    for ch_key, ch_data in channels.items():
        entry = {
            "label": ch_data.get("label", ch_key.replace("_", " ").title()),
            "pct": ch_data.get("pct") or 0,
            "color": ch_data.get("color", BLUE),
            "dollar": 0,
            "projected_apps": 0,
            "projected_hires": 0,
            "cpa": 0,
        }
        # Match with budget engine data
        ba_match = ba_channel_alloc.get(ch_key)
        if not ba_match:
            ch_label_lower = (ch_data.get("label") or "").lower()
            for ba_key, ba_val in ba_channel_alloc.items():
                if isinstance(ba_val, dict):
                    ba_label = ba_val.get("label", ba_key).lower()
                    if ba_label == ch_label_lower or ba_key.lower() == ch_key.lower():
                        ba_match = ba_val
                        break
        if ba_match and isinstance(ba_match, dict):
            entry["dollar"] = ba_match.get("dollar_amount") or 0
            entry["projected_apps"] = ba_match.get("projected_applications") or 0
            entry["projected_hires"] = ba_match.get("projected_hires") or 0
            entry["cpa"] = ba_match.get("cpa") or 0
        if ch_key in _reconciled_pct:
            entry["pct"] = _reconciled_pct[ch_key]
        # Fallback: compute dollar from percentage if budget engine didn't provide it
        if entry["dollar"] == 0 and ba_total_budget > 0 and entry["pct"] > 0:
            entry["dollar"] = ba_total_budget * entry["pct"] / 100

        display_channels.append(entry)

    # Sort by dollar amount (descending), then by percentage
    display_channels.sort(key=lambda c: (c["dollar"], c["pct"]), reverse=True)

    # Table header row
    header_y = table_top + Inches(0.35)
    row_h = Inches(0.38)
    col_widths = [
        Inches(3.0),
        Inches(1.8),
        Inches(2.2),
        Inches(1.8),
        Inches(1.8),
        Inches(1.5),
    ]
    col_headers = [
        "Channel",
        "Allocation %",
        "Investment",
        "Proj. Apps",
        "Proj. Hires",
        "CPA",
    ]
    col_aligns = [
        PP_ALIGN.LEFT,
        PP_ALIGN.CENTER,
        PP_ALIGN.CENTER,
        PP_ALIGN.CENTER,
        PP_ALIGN.CENTER,
        PP_ALIGN.CENTER,
    ]

    # Header background
    _add_filled_rect(slide, table_left, header_y, table_w, row_h, NAVY)

    cx = table_left
    for ci, (header, cw) in enumerate(zip(col_headers, col_widths)):
        _add_textbox(
            slide,
            cx + Inches(0.1),
            header_y,
            cw - Inches(0.1),
            row_h,
            text=header,
            font_size=9,
            bold=True,
            color=WHITE,
            alignment=col_aligns[ci],
            anchor=MSO_ANCHOR.MIDDLE,
        )
        cx += cw

    # S6 fix: this table used to hard-cap at 6 rows (``min(len(display_channels), 6)``)
    # while the Total row below summed over ALL channels -- on an 8-channel plan
    # that silently dropped the 2 smallest channels from the visible rows, so the
    # 6 shown rows never footed to the printed Total (confirmed on the Pratt & Whitney
    # NZ deck: 6 rows summed to NZ$142,499 / 96% against a printed NZ$150,000 / 101%
    # Total). Apply the same fit-aware cap the budget pie chart already uses
    # (top-(N-1) by dollar amount, tail rolled into one labeled "Other" row) so the
    # VISIBLE rows -- including the rollup -- always sum exactly to the Total.
    _MAX_VISIBLE_ROWS = 6
    if len(display_channels) > _MAX_VISIBLE_ROWS:
        _visible = display_channels[: _MAX_VISIBLE_ROWS - 1]
        _tail = display_channels[_MAX_VISIBLE_ROWS - 1 :]
        _tail_dollar = sum(c["dollar"] for c in _tail)
        _tail_apps = sum(int(c["projected_apps"] or 0) for c in _tail)
        _tail_hires = sum(int(c["projected_hires"] or 0) for c in _tail)
        _tail_pct = sum(c["pct"] for c in _tail)
        _tail_cpa = (_tail_dollar / _tail_apps) if _tail_apps > 0 else 0
        rollup_entry = {
            "label": f"+{len(_tail)} smaller channels",
            "pct": _tail_pct,
            "color": MUTED_TEXT,
            "dollar": _tail_dollar,
            "projected_apps": _tail_apps,
            "projected_hires": _tail_hires,
            "cpa": _tail_cpa,
            "_is_rollup": True,
        }
        rows_to_render = _visible + [rollup_entry]
    else:
        rows_to_render = list(display_channels)

    max_rows = len(rows_to_render)
    for ri, ch in enumerate(rows_to_render):
        row_y = header_y + row_h + ri * row_h
        row_bg = WHITE if ri % 2 == 0 else RGBColor(0xF5, 0xF5, 0xF3)
        _add_filled_rect(slide, table_left, row_y, table_w, row_h, row_bg)

        # Color indicator dot + Channel name
        dot_size = Inches(0.12)
        _add_oval(
            slide,
            table_left + Inches(0.12),
            row_y + (row_h - dot_size) / 2,
            dot_size,
            dot_size,
            ch["color"],
        )

        row_values = [
            ch["label"],
            f"{ch['pct']}%",
            f"{_cur}{ch['dollar']:,.0f}" if ch["dollar"] > 0 else "--",
            f"{int(ch['projected_apps']):,}" if ch["projected_apps"] > 0 else "--",
            f"{int(ch['projected_hires']):,}" if ch["projected_hires"] > 0 else "--",
            f"{_cur}{ch['cpa']:,.0f}" if ch["cpa"] > 0 else "--",
        ]

        cx = table_left
        for ci, (val, cw) in enumerate(zip(row_values, col_widths)):
            left_pad = Inches(0.3) if ci == 0 else Inches(0.1)
            _add_textbox(
                slide,
                cx + left_pad,
                row_y,
                cw - left_pad,
                row_h,
                text=val,
                font_size=9,
                bold=(ci == 0 and not ch.get("_is_rollup")),
                italic=bool(ch.get("_is_rollup")),
                color=DARK_TEXT,
                alignment=col_aligns[ci],
                anchor=MSO_ANCHOR.MIDDLE,
            )
            cx += cw

    # ---- TOTALS ROW ----
    # Executive decks expect a footed table. Totals are summed over ALL display
    # channels (not just the visible/rolled-up rows) -- but now the visible rows
    # (including the "+N smaller channels" rollup, if present) sum to exactly
    # this same total, so the table always foots.
    total_pct = sum(c["pct"] for c in display_channels)
    total_dollar = sum(c["dollar"] for c in display_channels)
    total_apps = sum(int(c["projected_apps"] or 0) for c in display_channels)
    total_hires = sum(int(c["projected_hires"] or 0) for c in display_channels)
    # Blended CPA = total spend / total apps (more meaningful than a CPA sum).
    blended_cpa = (total_dollar / total_apps) if total_apps > 0 else 0
    totals_y = header_y + row_h + max_rows * row_h
    _add_filled_rect(slide, table_left, totals_y, table_w, row_h, NAVY)
    totals_values = [
        "Total",
        f"{round(total_pct)}%",
        f"{_cur}{total_dollar:,.0f}" if total_dollar > 0 else "--",
        f"{total_apps:,}" if total_apps > 0 else "--",
        f"{total_hires:,}" if total_hires > 0 else "--",
        f"{_cur}{blended_cpa:,.0f}" if blended_cpa > 0 else "--",
    ]
    cx = table_left
    for ci, (val, cw) in enumerate(zip(totals_values, col_widths)):
        left_pad = Inches(0.3) if ci == 0 else Inches(0.1)
        _add_textbox(
            slide,
            cx + left_pad,
            totals_y,
            cw - left_pad,
            row_h,
            text=val,
            font_size=9,
            bold=True,
            color=WHITE,
            alignment=col_aligns[ci],
            anchor=MSO_ANCHOR.MIDDLE,
        )
        cx += cw

    # ---- FUNNEL STRIP (S93: funnel-calibration model) ----
    # HARD INVARIANT: every number here is read verbatim from
    # metadata.funnel.totals -- clicks/applications/hires match the table
    # above exactly; qualified/interviews are ADDED explanatory stages, not
    # a recomputation. Purely additive context so hires never sit directly
    # next to raw applications with no visible intermediate funnel.
    # max_rows is hard-capped at _MAX_VISIBLE_ROWS (rollup absorbs any
    # overflow into one extra row) -- so totals_y + row_h is a FIXED
    # worst-case value regardless of channel count, and the vertical band
    # between it and the footer rule (y=7.12in, see _add_footer) never
    # grows. To fit the new strip in that fixed band without colliding with
    # the footer, the combined (gap + strip + gap + insight box) height is
    # kept to the SAME 0.78in total the insight box alone used to occupy
    # (0.16in gap + 0.62in box) -- just re-subdivided, so the insight box's
    # bottom edge lands at the exact same y it always has.
    funnel_meta = ba_metadata.get("funnel") if isinstance(ba_metadata, dict) else None
    funnel_strip_h = Inches(0)
    strip_gap_before = Inches(0.05)
    strip_gap_after = Inches(0.03)
    strip_top = totals_y + row_h + strip_gap_before
    if isinstance(funnel_meta, dict):
        f_totals = funnel_meta.get("totals") or {}
        f_raw = f_totals.get("raw_apps")
        f_qual = f_totals.get("qualified_apps")
        f_int = f_totals.get("interviews")
        f_hires = f_totals.get("hires")
        f_clicks = ba_total_proj.get("clicks") or 0
        if all(v is not None for v in (f_raw, f_qual, f_int, f_hires)):
            funnel_strip_h = Inches(0.20)
            funnel_text = (
                f"Funnel: {int(f_clicks):,} clicks → {int(f_raw):,} applications "
                f"→ {int(f_qual):,} qualified → {int(f_int):,} interviews "
                f"→ {int(f_hires):,} hires"
            )
            _add_textbox(
                slide,
                table_left,
                strip_top,
                table_w,
                funnel_strip_h,
                text=funnel_text,
                font_size=9,
                italic=True,
                color=MUTED_TEXT,
                alignment=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.MIDDLE,
            )

    # ---- ROI INSIGHT CALLOUT (reference-deck dark-indigo takeaway band) ----
    # Positioned dynamically below the totals row (and the funnel strip, if
    # present) so it never overlaps the table (the old fixed y=6.05 collided
    # with the totals row when the table was tall) OR the footer rule below.
    if funnel_strip_h:
        insight_h = Inches(0.50)
        insight_top = strip_top + funnel_strip_h + strip_gap_after
    else:
        insight_h = Inches(0.62)
        insight_top = totals_y + row_h + Inches(0.16)
    _add_rounded_rect(slide, Inches(0.55), insight_top, Inches(12.2), insight_h, NAVY)
    _add_filled_rect(slide, Inches(0.55), insight_top, Inches(0.06), insight_h, TEAL)

    # Build insight text
    avg_cpa = ba_total_proj.get("cost_per_application") or 0
    # copy:both#2: derive the blended CPH the same way the Executive Summary
    # slide does (_compute_blended_cph), instead of an independent raw
    # ``total_projected.cost_per_hire`` read -- the two used to disagree
    # ("$5,263.16/hire" on Slide 2 vs. an unexplained "$5,250" here).
    avg_cph, _ = _compute_blended_cph(budget_alloc)

    if avg_cpa and avg_cpa > 0 and proj_hires and proj_hires > 0:
        insight_text = (
            f"Budget engine projects {_cur}{avg_cpa:,.0f} average CPA across all channels"
        )
        if avg_cph and avg_cph > 0:
            insight_text += f", with {_fmt.fmt_money(avg_cph)} average cost-per-hire"
        insight_text += (
            f". At {int(proj_hires):,} projected hires, "
            f"{client}'s investment yields strong programmatic ROI "
            f"through ML-driven bid optimization."
        )
    elif ba_total_budget > 0:
        insight_text = (
            f"{client}'s {_cur}{ba_total_budget:,.0f} investment is distributed across "
            f"{n_channels} channels using Nova AI Suite's programmatic optimization engine, "
            f"maximizing reach and conversion through real-time bid management."
        )
    else:
        insight_text = (
            f"Nova AI Suite's programmatic engine distributes {client}'s budget across "
            f"{n_channels} optimized channels with ML-driven bid management, "
            f"ensuring maximum ROI through continuous performance optimization."
        )

    _add_textbox(
        slide,
        Inches(0.85),
        insight_top,
        Inches(11.6),
        insight_h,
        text=_trunc_word(insight_text, 320),
        font_size=10,
        color=WHITE,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # Footer
    _add_footer(slide, today)


def _embed_pie_chart_on_budget_slide(prs: Presentation, data: Dict) -> None:
    """Embed a small budget pie chart onto the last slide (budget allocation).

    S48: Instead of a separate pie chart slide, this places a compact
    pie chart image in the bottom-right area of the budget allocation slide.
    Called immediately after _build_slide_budget_allocation in generate_pptx.
    """
    if not prs.slides or len(prs.slides) == 0:
        return

    channels = _selected_channels(data)
    if not channels:
        return

    budget_alloc = data.get("_budget_allocation", {})
    ba_channel_alloc = (
        budget_alloc.get("channel_allocations", {})
        if isinstance(budget_alloc, dict)
        else {}
    )

    labels: List[str] = []
    sizes: List[float] = []

    # S6 fix: reconcile ALL channels' percentages together (largest-remainder
    # rounding) instead of rounding each independently, so this pie's legend
    # percentages sum to exactly 100 instead of drifting to 101%/99%.
    _reconciled_pct = _reconcile_channel_percentages(channels, ba_channel_alloc)

    for ch_key, ch_data in channels.items():
        label = ch_data.get("label", ch_key.replace("_", " ").title())
        pct = _reconciled_pct.get(ch_key, ch_data.get("pct") or 0)

        if pct > 0:
            labels.append(label)
            sizes.append(pct)

    if not labels:
        return

    # Generate a compact pie chart (smaller than the standalone version)
    try:
        fig, ax = plt.subplots(figsize=(3.5, 2.8), dpi=150)
        fig.patch.set_facecolor("#FFFCF9")

        colors = _CHART_COLORS[: len(labels)]
        while len(colors) < len(labels):
            colors.append(_CHART_COLORS[len(colors) % len(_CHART_COLORS)])

        wedges, texts, autotexts = ax.pie(
            sizes,
            labels=None,
            autopct=lambda pct: f"{pct:.0f}%" if pct >= 5 else "",
            startangle=90,
            colors=colors,
            pctdistance=0.75,
            wedgeprops={"linewidth": 1, "edgecolor": "white"},
        )

        for autotext in autotexts:
            autotext.set_fontsize(7)
            autotext.set_fontweight("bold")
            autotext.set_color("white")

        legend_labels = [f"{lbl} ({sz:.0f}%)" for lbl, sz in zip(labels, sizes)]
        ax.legend(
            wedges,
            legend_labels,
            loc="center left",
            bbox_to_anchor=(1.0, 0.5),
            fontsize=6,
            frameon=False,
        )

        plt.tight_layout()
        buf = io.BytesIO()
        fig.savefig(
            buf, format="png", bbox_inches="tight", facecolor=fig.get_facecolor()
        )
        plt.close(fig)
        buf.seek(0)
        chart_bytes = buf.getvalue()
    except Exception as exc:
        logger.warning("Embedded pie chart generation failed: %s", exc)
        return

    # Place the chart on the last slide (budget allocation) -- bottom-right
    budget_slide = prs.slides[len(prs.slides) - 1]
    chart_stream = io.BytesIO(chart_bytes)
    budget_slide.shapes.add_picture(
        chart_stream,
        Inches(8.8),  # right side
        Inches(5.5),  # below table rows, above footer
        Inches(4.2),  # width
        Inches(1.9),  # height
    )


def _cmp_status(
    client_val: Any, industry_val: Any, higher_is_better: bool = True, tie_band: float = 0.05
) -> str:
    """Honest 3-state benchmark comparison: ``'beating'`` / ``'on_par'`` /
    ``'trailing'``. A client value within +/-``tie_band`` (default 5%) of the
    benchmark -- including an exact tie -- is ``'on_par'``, never flagged as
    beating (a shipped bundle previously hardcoded ``is_better: True`` for
    rows with no real comparison at all)."""
    try:
        c = float(client_val)
        i = float(industry_val)
    except (TypeError, ValueError):
        return "on_par"
    if i == 0:
        if c == 0:
            return "on_par"
        is_positive_direction = c > 0
        return (
            "beating"
            if is_positive_direction == higher_is_better
            else "trailing"
        )
    diff_pct = (c - i) / abs(i)
    if abs(diff_pct) <= tie_band:
        return "on_par"
    if higher_is_better:
        return "beating" if diff_pct > 0 else "trailing"
    return "beating" if diff_pct < 0 else "trailing"


def _cmp_status_in_range(
    value: float, low: float, high: float, higher_is_better: bool = True
) -> str:
    """3-state status for a value against a benchmark RANGE (not a single
    point): inside ``[low, high]`` is always ``'on_par'`` ("within range"),
    never ``'beating'`` even when the value sits at the favorable end.
    Outside the range, direction is governed by ``higher_is_better``
    (default True -- e.g. more markets covered than the healthy range is
    'beating', fewer is 'trailing'; pass False for metrics like CPA where
    lower is better)."""
    if low > high:
        low, high = high, low
    if low <= value <= high:
        return "on_par"
    above = value > high
    return "beating" if (above == higher_is_better) else "trailing"


_CMP_STATUS_STYLE = {
    "beating": (GREEN, "▲", "Beating benchmark"),
    "trailing": (AMBER, "▼", "Trailing benchmark"),
    "on_par": (MUTED_TEXT, "—", "On par / within range"),
    "none": (MUTED_TEXT, "", ""),
}


# ===================================================================
# SLIDE 7 - Side-by-Side Comparison Panel + Implementation Timeline
# ===================================================================


def _build_slide_comparison_timeline(prs: Presentation, data: Dict):
    """Build comparison panel (Client Plan vs Industry Average) and implementation timeline."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    client = data.get("client_name", "Client")
    industry = data.get("industry", "general_entry_level")
    industry_label = data.get("industry_label", industry.replace("_", " ").title())
    channels = _selected_channels(data)
    budget = data.get("budget", "TBD")
    locations = data.get("locations") or []
    roles = data.get("roles") or []
    today = datetime.date.today().strftime("%B %d, %Y")

    # Pull synthesized + budget allocation data (from pipeline)
    synthesized = data.get("_synthesized", {})
    budget_alloc = data.get("_budget_allocation", {})

    # Off-white background
    _add_filled_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, OFF_WHITE)

    # Top band
    _add_top_band(slide, "PLAN COMPARISON & IMPLEMENTATION", today)

    # Action title
    action_text = (
        f"{client}'s optimized media plan vs. {industry_label} industry averages "
        f"with phased implementation roadmap"
    )
    _add_textbox(
        slide,
        Inches(0.55),
        Inches(0.92),
        Inches(12.2),
        Inches(0.45),
        text=action_text,
        font_size=15,
        bold=True,
        color=NAVY,
    )

    # ---- SIDE-BY-SIDE COMPARISON ----
    comp_top = Inches(1.55)
    panel_h = Inches(2.95)
    panel_w = Inches(5.9)
    panel_gap = Inches(0.4)
    left_panel_x = Inches(0.55)
    right_panel_x = left_panel_x + panel_w + panel_gap

    ind_benchmarks = _get_industry_comparison(industry, data)
    n_channels = len(channels)
    n_locations = len(locations)

    # visual:atria#1 / strategy:atria#1 / consistency:atria#1 / copy:both#1
    # (4-critical cluster): ``channels`` (from ``_selected_channels``) only
    # carries the STATIC INDUSTRY_ALLOC_PROFILES percentages -- slides 5/6
    # override those with the final post-reweight ``_budget_allocation``
    # figures before reading them (see ``_build_slide_channel_strategy`` /
    # ``_build_slide_budget_allocation``), but this slide never did, so its
    # "Programmatic Allocation" (and every other plan-column row derived
    # from ``channels[*]["pct"]``) showed the pre-reweight profile split
    # (e.g. 21%/25%) instead of the plan's actual committed spend (28%/26%).
    # Apply the SAME reconciliation used on slides 5/6 so every plan-column
    # value here reads the final allocation data.
    _cmp_ba_channel_alloc = (
        budget_alloc.get("channel_allocations", {}) if budget_alloc else {}
    )
    if _cmp_ba_channel_alloc:
        _cmp_reconciled_pct = _reconcile_channel_percentages(
            channels, _cmp_ba_channel_alloc
        )
        for ch_key, ch_data in channels.items():
            if ch_key in _cmp_reconciled_pct:
                ch_data["pct"] = _cmp_reconciled_pct[ch_key]

    # Calculate client metrics
    sorted_ch = sorted(channels.values(), key=lambda c: c["pct"], reverse=True)
    programmatic_pct = 0
    for ch in channels.values():
        if ch.get("category") == "Programmatic":
            programmatic_pct += ch["pct"]
    if programmatic_pct == 0:
        programmatic_pct = sorted_ch[0]["pct"] if sorted_ch else 30

    # Comparison metrics - build all candidates
    # Honest 3-state comparison ('beating' / 'on_par' / 'trailing') via
    # _cmp_status -- a +/-5% tie band means an exact or near-tie (e.g. 6
    # channels vs an industry average of 6) renders neutral, never a green
    # "beating" arrow.
    #
    # consistency:atria#3 / consistency:manpower#4: "Reach Multiplier" and
    # "Channel Diversity Score" used to be computed here from a made-up
    # formula (``1.0 + (n_channels - 4) * 0.15`` / ``min(10, n_channels *
    # 1.5)``) with no workbook source -- neither figure traces to anything
    # in ``_budget_allocation`` or the KB, so they were removed rather than
    # kept as unsourced client-facing statistics. Every remaining row here
    # traces to ``_selected_channels``/``_budget_allocation``/the industry
    # benchmark KB.
    _geo_status = _cmp_status_in_range(n_locations, 3, 5)
    all_comparison_rows = [
        {
            "metric": "Channels Selected",
            "client_val": str(n_channels),
            "industry_val": str(ind_benchmarks.get("avg_channels", 4)),
            "status": _cmp_status(n_channels, ind_benchmarks.get("avg_channels", 4)),
        },
        {
            "metric": "Programmatic Allocation",
            "client_val": f"{programmatic_pct}%",
            "industry_val": f"{ind_benchmarks.get('avg_budget_pct_programmatic', 30)}%",
            "status": _cmp_status(
                programmatic_pct, ind_benchmarks.get("avg_budget_pct_programmatic", 30)
            ),
        },
        {
            "metric": "Geographic Coverage",
            "client_val": f"{n_locations} market{'s' if n_locations != 1 else ''}",
            # A value INSIDE the 3-5 benchmark range is 'within range' -- it
            # is never rendered as "beating" just for landing inside it.
            "industry_val": "3-5 markets (within range)" if _geo_status == "on_par" else "3-5 markets",
            "status": _geo_status,
        },
    ]

    # Add budget-allocation-powered comparison rows if real data is available
    ba_total_proj_comp = budget_alloc.get("total_projected", {}) if budget_alloc else {}
    if not isinstance(ba_total_proj_comp, dict):
        ba_total_proj_comp = {}
    ba_channel_alloc = (
        budget_alloc.get("channel_allocations", {}) if budget_alloc else {}
    )
    ba_metadata_comp = budget_alloc.get("metadata", {}) if budget_alloc else {}
    if not isinstance(ba_metadata_comp, dict):
        ba_metadata_comp = {}
    ba_total_budget = ba_metadata_comp.get("total_budget") or 0

    proj_hires = 0
    comp_cph = 0.0
    if ba_total_proj_comp:
        proj_cpa = ba_total_proj_comp.get("cost_per_application") or 0
        # S48 FIX: per-channel sum for consistency
        if isinstance(ba_channel_alloc, dict) and ba_channel_alloc:
            proj_hires = sum(
                int(ch.get("projected_hires") or 0) for ch in ba_channel_alloc.values()
            )
        else:
            proj_hires = 0
        if proj_hires == 0:
            proj_hires = ba_total_proj_comp.get("hires") or 0
        proj_apps = ba_total_proj_comp.get("applications") or 0
        comp_cph = (
            round(ba_total_budget / proj_hires, 2)
            if proj_hires > 0 and ba_total_budget > 0
            else (ba_total_proj_comp.get("cost_per_hire") or 0)
        )

        # Get industry benchmark CPA for comparison -- a low/high RANGE, not
        # just an average, so a value inside the range renders 'within
        # range' (on_par) rather than 'beating' the moment it clears the
        # midpoint.
        bench = _get_benchmarks(industry, data)
        cpa_str = bench.get("cpa", "$25")
        try:
            cpa_nums = [float(x) for x in re.findall(r"[\d.]+", cpa_str.replace(",", ""))]
        except Exception:
            cpa_nums = []
        if len(cpa_nums) >= 2:
            cpa_low, cpa_high = min(cpa_nums), max(cpa_nums)
        elif len(cpa_nums) == 1:
            cpa_low = cpa_high = cpa_nums[0]
        else:
            cpa_low = cpa_high = 25.0

        if proj_cpa and proj_cpa > 0:
            # S3: cpa_str may still be a fixed US-benchmark constant (bare $)
            # while client_val is already in the plan's local currency -- never
            # show them side-by-side with no marker. Mark the industry-side
            # figure itself (not just a label) when it's a USD benchmark on a
            # non-USD plan.
            _cpa_industry_val = cpa_str
            if (
                bench.get("cpa_is_usd_benchmark", True)
                and _get_active_currency() != "USD"
            ):
                _cpa_industry_val = _mark_usd(cpa_str)
            if cpa_low == cpa_high:
                # Point estimate (not a range): lower CPA is better.
                _cpa_status = _cmp_status(proj_cpa, cpa_low, higher_is_better=False)
            elif cpa_low <= proj_cpa <= cpa_high:
                _cpa_status = "on_par"  # within range -- never "beating"
            else:
                _cpa_status = "beating" if proj_cpa < cpa_low else "trailing"
            all_comparison_rows.append(
                {
                    "metric": "Projected CPA",
                    "client_val": _fmt_currency(proj_cpa),
                    "industry_val": _cpa_industry_val,
                    "status": _cpa_status,
                }
            )
        if proj_hires and proj_hires > 0:
            # No real industry-average "hires" benchmark exists to compare
            # against -- render a neutral em-dash with NO arrow/badge rather
            # than a fabricated "beating" claim.
            all_comparison_rows.append(
                {
                    "metric": "Projected Hires",
                    "client_val": f"{int(proj_hires):,}",
                    "industry_val": "—",
                    "status": "none",
                }
            )
        if ba_total_budget and ba_total_budget > 0:
            # Same: "Varies" is not a real benchmark to compare against.
            all_comparison_rows.append(
                {
                    "metric": "Total Investment",
                    "client_val": _fmt_currency(ba_total_budget),
                    "industry_val": "Varies",
                    "status": "none",
                }
            )

    # ---- Honest goal-gap row (display_format.parse_hire_goal / goal_gap) ----
    # When the client stated a hiring GOAL and this plan's projection falls
    # short of it, add an honest gap callout -- framed as a scaling roadmap,
    # never as a failure.
    _hire_goal = _fmt.parse_hire_goal(data.get("hire_volume"))
    _goal_gap = _fmt.goal_gap(proj_hires, _hire_goal, comp_cph)
    if _goal_gap:
        all_comparison_rows.append(
            {
                "metric": "vs. Client Goal",
                "client_val": f"{_goal_gap['projected']:,} ({_goal_gap['pct_of_goal']:.0f}%)",
                "industry_val": f"{_goal_gap['goal']:,} hires goal",
                "status": "trailing",
            }
        )

    # Honest ordering: beating benchmark first, then on-par/within-range,
    # then trailing/no-comparison last. NOTE: this deliberately does NOT
    # relabel trailing metrics as "(Target)" and flip them to a green
    # "beating" arrow -- that used to cosmetically launder every trailing
    # metric into a fabricated win whenever trailing rows were the majority.
    _status_rank = {"beating": 0, "on_par": 1, "trailing": 2, "none": 3}
    comparison_rows = sorted(
        all_comparison_rows, key=lambda r: _status_rank.get(r["status"], 3)
    )
    # strategy:manpower#3: the client's own hiring GOAL is not an industry
    # statistic -- putting it in the "Industry Average" column (as this row
    # used to) is a category error that compares the plan against itself.
    # Pull it out of the two-panel benchmark table entirely; it always keeps
    # a slot (rendered as its own distinctly-labeled band below the panels,
    # see below) rather than risking getting crowded out by the row cap.
    _goal_row = None
    if _goal_gap:
        _goal_row = next(
            (r for r in comparison_rows if r["metric"] == "vs. Client Goal"), None
        )
        comparison_rows = [r for r in comparison_rows if r is not _goal_row]
    comparison_rows = comparison_rows[:4] if _goal_row else comparison_rows[:5]

    # ==== LEFT PANEL: Client Plan ====
    _add_rounded_rect(slide, left_panel_x, comp_top, panel_w, panel_h, WHITE)
    # Header bar
    _add_filled_rect(slide, left_panel_x, comp_top, panel_w, Inches(0.45), NAVY)
    _add_textbox(
        slide,
        left_panel_x + Inches(0.2),
        comp_top + Inches(0.05),
        panel_w - Inches(0.4),
        Inches(0.35),
        text=f"\u2b22  {client}'s Plan",
        font_size=12,
        bold=True,
        color=WHITE,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    row_h_comp = Inches(0.45)
    for ri, row in enumerate(comparison_rows):
        ry = comp_top + Inches(0.5) + ri * row_h_comp
        bg = WHITE if ri % 2 == 0 else RGBColor(0xF8, 0xF6, 0xF3)
        _add_filled_rect(
            slide,
            left_panel_x + Inches(0.05),
            ry,
            panel_w - Inches(0.1),
            row_h_comp,
            bg,
        )

        # Metric label
        _add_textbox(
            slide,
            left_panel_x + Inches(0.2),
            ry,
            Inches(2.8),
            row_h_comp,
            text=row["metric"],
            font_size=9,
            bold=True,
            color=DARK_TEXT,
            anchor=MSO_ANCHOR.MIDDLE,
        )

        # Value with status indicator -- 3-state (beating / on_par /
        # trailing / none), a tie or "no real benchmark" row renders a
        # neutral em-dash with no colored arrow at all.
        status_color, indicator, _ = _CMP_STATUS_STYLE.get(
            row.get("status", "on_par"), _CMP_STATUS_STYLE["on_par"]
        )

        val_box, val_tf = _add_textbox(
            slide,
            left_panel_x + Inches(3.2),
            ry,
            Inches(2.5),
            row_h_comp,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        p = val_tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r1 = p.add_run()
        r1.text = row["client_val"]
        _set_font(r1, size=12, bold=True, color=NAVY)
        r2 = p.add_run()
        r2.text = f"  {indicator}"
        _set_font(r2, size=10, bold=True, color=status_color)

    # ==== RIGHT PANEL: Industry Average ====
    _add_rounded_rect(slide, right_panel_x, comp_top, panel_w, panel_h, WHITE)
    _add_filled_rect(slide, right_panel_x, comp_top, panel_w, Inches(0.45), MUTED_TEXT)
    _add_textbox(
        slide,
        right_panel_x + Inches(0.2),
        comp_top + Inches(0.05),
        panel_w - Inches(0.4),
        Inches(0.35),
        text=f"\u25cb  {industry_label} Average",
        font_size=12,
        bold=True,
        color=WHITE,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    for ri, row in enumerate(comparison_rows):
        ry = comp_top + Inches(0.5) + ri * row_h_comp
        bg = WHITE if ri % 2 == 0 else RGBColor(0xF8, 0xF6, 0xF3)
        _add_filled_rect(
            slide,
            right_panel_x + Inches(0.05),
            ry,
            panel_w - Inches(0.1),
            row_h_comp,
            bg,
        )

        _add_textbox(
            slide,
            right_panel_x + Inches(0.2),
            ry,
            Inches(2.8),
            row_h_comp,
            text=row["metric"],
            font_size=9,
            bold=True,
            color=DARK_TEXT,
            anchor=MSO_ANCHOR.MIDDLE,
        )

        _add_textbox(
            slide,
            right_panel_x + Inches(3.2),
            ry,
            Inches(2.5),
            row_h_comp,
            text=row["industry_val"],
            font_size=12,
            bold=True,
            color=MUTED_TEXT,
            alignment=PP_ALIGN.RIGHT,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    # ---- CLIENT GOAL band (own row, never inside the industry-average
    # column -- strategy:manpower#3) ----
    if _goal_row and _goal_gap:
        _goal_band_y = comp_top + Inches(0.5) + len(comparison_rows) * row_h_comp
        _goal_band_w = (right_panel_x + panel_w) - left_panel_x
        _add_filled_rect(
            slide,
            left_panel_x + Inches(0.05),
            _goal_band_y,
            _goal_band_w - Inches(0.1),
            row_h_comp,
            LIGHT_AMBER,
        )
        _gb_box, _gb_tf = _add_textbox(
            slide,
            left_panel_x + Inches(0.2),
            _goal_band_y,
            _goal_band_w - Inches(0.4),
            row_h_comp,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        _gb_p = _gb_tf.paragraphs[0]
        _gb_r1 = _gb_p.add_run()
        _gb_r1.text = "CLIENT GOAL (not an industry benchmark):  "
        _set_font(_gb_r1, size=9, bold=True, color=NAVY)
        _gb_r2 = _gb_p.add_run()
        _gb_r2.text = (
            f"{_goal_gap['goal']:,} hires target  —  this plan projects "
            f"{_goal_gap['projected']:,} ({_goal_gap['pct_of_goal']:.0f}%)"
        )
        _set_font(_gb_r2, size=9, bold=False, color=DARK_TEXT)

    # ---- Legend ----
    legend_y = comp_top + panel_h + Inches(0.04)
    leg_box, leg_tf = _add_textbox(
        slide,
        Inches(0.55),
        legend_y,
        Inches(6),
        Inches(0.22),
    )
    p = leg_tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "\u25b2 "
    _set_font(r1, size=8, bold=True, color=GREEN)
    r2 = p.add_run()
    r2.text = "Beating benchmark    "
    _set_font(r2, size=8, color=MUTED_TEXT)
    r3 = p.add_run()
    r3.text = "\u25bc "
    _set_font(r3, size=8, bold=True, color=AMBER)
    r4 = p.add_run()
    r4.text = "Trailing benchmark    "
    _set_font(r4, size=8, color=MUTED_TEXT)
    r5 = p.add_run()
    r5.text = "\u2014 "
    _set_font(r5, size=8, bold=True, color=MUTED_TEXT)
    r6 = p.add_run()
    r6.text = "On par / within range"
    _set_font(r6, size=8, color=MUTED_TEXT)

    # ==== IMPLEMENTATION TIMELINE (bottom) ====
    timeline_top = Inches(4.98)  # was 4.8 — header overlapped the legend above

    _add_textbox(
        slide,
        Inches(0.55),
        timeline_top,
        Inches(12.2),
        Inches(0.32),
        text="IMPLEMENTATION TIMELINE",
        font_size=11,
        bold=True,
        color=NAVY,
    )
    _add_filled_rect(
        slide, Inches(0.55), timeline_top + Inches(0.3), Inches(2.2), Inches(0.03), TEAL
    )

    # Build timeline phases based on actual campaign_weeks from input
    cw = data.get("campaign_weeks", 12)
    if cw <= 12:
        p2_end = min(6, cw)
        p3_start = min(7, cw)
        p3_end = cw
        phases = [
            {
                "phase": "PHASE 1",
                "weeks": "Weeks 1-2",
                "title": "Launch & Calibrate",
                "bullets": [
                    "Campaign setup & publisher activation",
                    "Baseline measurement & tracking",
                    "Attribution configuration",
                ],
                "color": BLUE,
                "accent_bg": LIGHT_BLUE,
            },
            {
                "phase": "PHASE 2",
                "weeks": f"Weeks 3-{p2_end}",
                "title": "Optimize & Scale",
                "bullets": [
                    "ML bid optimization active",
                    "A/B test creative & targeting",
                    "Scale top performers",
                ],
                "color": GREEN,
                "accent_bg": LIGHT_GREEN,
            },
            {
                "phase": "PHASE 3",
                "weeks": f"Weeks {p3_start}-{p3_end}",
                "title": "Maximize & Report",
                "bullets": [
                    "Full CPQA optimization",
                    "ROI analysis & reallocation",
                    "Performance review",
                ],
                "color": NAVY,
                "accent_bg": RGBColor(0xE8, 0xED, 0xF4),
            },
        ]
    elif cw <= 26:
        phases = [
            {
                "phase": "PHASE 1",
                "weeks": "Weeks 1-3",
                "title": "Launch & Calibrate",
                "bullets": [
                    "Campaign setup & publisher activation",
                    "Baseline measurement & tracking",
                    "Attribution configuration",
                ],
                "color": BLUE,
                "accent_bg": LIGHT_BLUE,
            },
            {
                "phase": "PHASE 2",
                "weeks": f"Weeks 4-{cw // 2}",
                "title": "Optimize & Scale",
                "bullets": [
                    "ML bid optimization active",
                    "A/B test creative & targeting",
                    "Scale top performers",
                ],
                "color": GREEN,
                "accent_bg": LIGHT_GREEN,
            },
            {
                "phase": "PHASE 3",
                "weeks": f"Weeks {cw // 2 + 1}-{cw}",
                "title": "Maximize & Report",
                "bullets": [
                    "Full CPQA optimization",
                    "ROI analysis & reallocation",
                    "Quarterly performance review",
                ],
                "color": NAVY,
                "accent_bg": RGBColor(0xE8, 0xED, 0xF4),
            },
        ]
    else:
        phases = [
            {
                "phase": "PHASE 1",
                "weeks": "Weeks 1-4",
                "title": "Launch & Calibrate",
                "bullets": [
                    "Campaign setup & publisher activation",
                    "Baseline measurement & tracking",
                    "Attribution configuration",
                ],
                "color": BLUE,
                "accent_bg": LIGHT_BLUE,
            },
            {
                "phase": "PHASE 2",
                "weeks": f"Weeks 5-{cw // 3}",
                "title": "Optimize & Scale",
                "bullets": [
                    "ML bid optimization active",
                    "A/B test creative & targeting",
                    "Scale top performers",
                ],
                "color": GREEN,
                "accent_bg": LIGHT_GREEN,
            },
            {
                "phase": "PHASE 3",
                "weeks": f"Weeks {cw // 3 + 1}-{cw}",
                "title": "Maximize & Report",
                "bullets": [
                    "Full CPQA optimization",
                    "ROI analysis & reallocation",
                    "Quarterly performance review",
                ],
                "color": NAVY,
                "accent_bg": RGBColor(0xE8, 0xED, 0xF4),
            },
        ]

    # copy:both#2: interpolate this plan's own facts (client name, top
    # channels) into the otherwise-identical phase action items.
    _interpolate_timeline_bullets(phases, data, channels)

    phase_w = Inches(3.85)
    phase_gap = Inches(0.25)
    phase_top = timeline_top + Inches(0.45)
    phase_h = Inches(1.65)

    for i, ph in enumerate(phases):
        px = Inches(0.55) + i * (phase_w + phase_gap)

        # Phase card
        _add_rounded_rect(slide, px, phase_top, phase_w, phase_h, WHITE)

        # Top accent bar
        _add_filled_rect(slide, px, phase_top, phase_w, Inches(0.05), ph["color"])

        # Phase number badge
        badge_w = Inches(0.9)
        badge_h = Inches(0.25)
        _add_rounded_rect(
            slide,
            px + Inches(0.12),
            phase_top + Inches(0.15),
            badge_w,
            badge_h,
            ph["accent_bg"],
        )
        _add_textbox(
            slide,
            px + Inches(0.12),
            phase_top + Inches(0.15),
            badge_w,
            badge_h,
            text=ph["phase"],
            font_size=7,
            bold=True,
            color=ph["color"],
            alignment=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )

        # Weeks
        _add_textbox(
            slide,
            px + Inches(1.1),
            phase_top + Inches(0.15),
            Inches(1.5),
            badge_h,
            text=ph["weeks"],
            font_size=8,
            color=MUTED_TEXT,
            anchor=MSO_ANCHOR.MIDDLE,
        )

        # Title
        _add_textbox(
            slide,
            px + Inches(0.12),
            phase_top + Inches(0.48),
            phase_w - Inches(0.24),
            Inches(0.25),
            text=ph["title"],
            font_size=11,
            bold=True,
            color=DARK_TEXT,
        )

        # Bullets
        bx, btf = _add_textbox(
            slide,
            px + Inches(0.12),
            phase_top + Inches(0.78),
            phase_w - Inches(0.24),
            Inches(0.85),
        )
        btf.paragraphs[0].space_before = Pt(0)
        btf.paragraphs[0].space_after = Pt(0)

        for j, bullet in enumerate(ph["bullets"]):
            if j == 0:
                bp = btf.paragraphs[0]
            else:
                bp = btf.add_paragraph()
            bp.space_before = Pt(1)
            bp.space_after = Pt(3)
            bp.alignment = PP_ALIGN.LEFT

            br = bp.add_run()
            br.text = "\u2713  "
            _set_font(br, size=8, bold=False, color=ph["color"])
            bt = bp.add_run()
            bt.text = bullet
            _set_font(bt, size=8, color=MUTED_TEXT)

    # Arrow connectors between phases
    for i in range(2):
        ax = (
            Inches(0.55) + (i + 1) * (phase_w + phase_gap) - phase_gap / 2 - Inches(0.1)
        )
        ay = phase_top + phase_h / 2 - Inches(0.1)
        _add_textbox(
            slide,
            ax,
            ay,
            Inches(0.2),
            Inches(0.2),
            text="\u25b6",
            font_size=12,
            bold=True,
            color=TEAL,
            alignment=PP_ALIGN.CENTER,
        )

    # Footer
    _add_footer(slide, today)


# ===================================================================
# SLIDE - Market & Workforce Analysis (NEW)
# ===================================================================


def _build_slide_market_analysis(prs: Presentation, data: Dict):
    """Build the Market & Workforce Analysis slide.

    Uses:
    - job_market_demand: market temperature, trends, macro-economic data
    - workforce_insights: Gen-Z trends, employer branding, research
    - salary_intelligence: salary ranges per role
    """
    try:
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        client = data.get("client_name", "Client")
        industry = data.get("industry", "general_entry_level")
        industry_label = data.get("industry_label", industry.replace("_", " ").title())
        roles = data.get("roles") or []
        today = datetime.date.today().strftime("%B %d, %Y")

        synthesized = data.get("_synthesized", {})
        if not isinstance(synthesized, dict):
            synthesized = {}
        job_market = synthesized.get("job_market_demand", {})
        if not isinstance(job_market, dict):
            job_market = {}
        workforce = synthesized.get("workforce_insights", {})
        if not isinstance(workforce, dict):
            workforce = {}
        salary_intel = synthesized.get("salary_intelligence", {})
        if not isinstance(salary_intel, dict):
            salary_intel = {}

        # Off-white background
        _add_filled_rect(
            slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, OFF_WHITE
        )

        # Top band
        _add_top_band(slide, "MARKET & WORKFORCE ANALYSIS", today)

        # Action title
        action_text = (
            f"Labor market intelligence and workforce trend analysis for "
            f"{client}'s {industry_label} hiring strategy"
        )
        _add_textbox(
            slide,
            Inches(0.55),
            Inches(0.92),
            Inches(12.2),
            Inches(0.5),
            text=action_text,
            font_size=15,
            bold=True,
            color=NAVY,
        )

        # ---- LEFT COLUMN: Market Demand by Role ----
        section_top = Inches(1.6)
        left_col_left = Inches(0.55)
        left_col_w = Inches(6.0)

        _add_textbox(
            slide,
            left_col_left,
            section_top,
            left_col_w,
            Inches(0.35),
            text="JOB MARKET DEMAND BY ROLE",
            font_size=11,
            bold=True,
            color=NAVY,
        )
        _add_filled_rect(
            slide,
            left_col_left,
            section_top + Inches(0.33),
            Inches(2.5),
            Inches(0.03),
            TEAL,
        )

        # Market demand table
        table_top = section_top + Inches(0.5)
        row_h = Inches(0.36)

        # Header
        _add_filled_rect(slide, left_col_left, table_top, left_col_w, row_h, NAVY)
        col_widths = [Inches(2.0), Inches(1.0), Inches(1.0), Inches(1.0), Inches(1.0)]
        col_headers = ["Role", "Postings", "Temp.", "Trend", "Competition"]
        cx = left_col_left
        for ci, (header, cw) in enumerate(zip(col_headers, col_widths)):
            _add_textbox(
                slide,
                cx + Inches(0.08),
                table_top,
                cw,
                row_h,
                text=header,
                font_size=8,
                bold=True,
                color=WHITE,
                anchor=MSO_ANCHOR.MIDDLE,
            )
            cx += cw

        # Data rows
        market_rows = []
        for role_name, role_data in list(job_market.items())[:6]:
            if not isinstance(role_data, dict):
                continue
            postings = role_data.get(
                "total_postings", role_data.get("posting_count") or 0
            )
            temp = role_data.get("market_temperature", "—")
            trend = role_data.get("trend_direction", "stable")
            comp_idx = role_data.get("competition_index") or 0
            market_rows.append(
                (
                    str(role_name)[:25],
                    (
                        f"{postings:,}"
                        if isinstance(postings, (int, float)) and postings > 0
                        else "—"
                    ),
                    temp.title() if temp else "—",
                    trend.title() if trend else "Stable",
                    (
                        f"{comp_idx:.2f}"
                        if isinstance(comp_idx, (int, float)) and comp_idx > 0
                        else "—"
                    ),
                )
            )

        if not market_rows:
            market_rows = [("Market data not available", "-", "-", "-", "-")]

        for ri, row_vals in enumerate(market_rows):
            ry = table_top + row_h * (ri + 1)
            bg = WHITE if ri % 2 == 0 else RGBColor(0xF8, 0xF6, 0xF3)
            _add_filled_rect(slide, left_col_left, ry, left_col_w, row_h, bg)
            cx = left_col_left
            for ci, (val, cw) in enumerate(zip(row_vals, col_widths)):
                # Color code temperature
                val_color = DARK_TEXT
                if ci == 2:  # Temperature column
                    if val.lower() == "hot":
                        val_color = RED_ACCENT
                    elif val.lower() == "warm":
                        val_color = AMBER
                    elif val.lower() == "cool":
                        val_color = BLUE
                    elif val.lower() == "cold":
                        val_color = MEDIUM_BLUE
                _add_textbox(
                    slide,
                    cx + Inches(0.08),
                    ry,
                    cw,
                    row_h,
                    text=val,
                    font_size=8,
                    bold=(ci == 0),
                    color=val_color,
                    anchor=MSO_ANCHOR.MIDDLE,
                )
                cx += cw

        # ---- Macro-Economic Context (below market table) ----
        macro_top = table_top + row_h * (len(market_rows) + 1) + Inches(0.25)
        _add_textbox(
            slide,
            left_col_left,
            macro_top,
            left_col_w,
            Inches(0.3),
            text="MACRO-ECONOMIC CONTEXT",
            font_size=10,
            bold=True,
            color=NAVY,
        )
        _add_filled_rect(
            slide,
            left_col_left,
            macro_top + Inches(0.28),
            Inches(2.0),
            Inches(0.03),
            TEAL,
        )

        # Extract macro data from first role's data
        macro_data = {}
        for _rk, _rv in job_market.items():
            if isinstance(_rv, dict) and _rv.get("macro_economic"):
                macro_data = _rv["macro_economic"]
                break

        # KB fallback: read fred_indicators.json directly from knowledge base
        if not macro_data:
            _kb = data.get("_knowledge_base", {})
            _fred_kb = _kb.get("fred_indicators", {}) if isinstance(_kb, dict) else {}
            _fred_data = (
                _fred_kb.get("data", _fred_kb) if isinstance(_fred_kb, dict) else {}
            )
            if isinstance(_fred_data, dict):
                _flat: Dict[str, Any] = {}
                for _fk, _fv in _fred_data.items():
                    if isinstance(_fv, dict) and "value" in _fv:
                        _flat[_fk] = _fv["value"]
                    elif isinstance(_fv, (int, float)):
                        _flat[_fk] = _fv
                if _flat:
                    macro_data = _flat

        macro_items = []
        if macro_data:
            unemp = macro_data.get("unemployment_rate")
            if unemp is not None:
                macro_items.append(
                    (
                        "Unemployment Rate",
                        f"{unemp}%" if isinstance(unemp, (int, float)) else str(unemp),
                    )
                )
            lfpr = macro_data.get("labor_force_participation")
            if lfpr is not None:
                macro_items.append(
                    (
                        "Labor Force Participation",
                        f"{lfpr}%" if isinstance(lfpr, (int, float)) else str(lfpr),
                    )
                )
            jolts = macro_data.get("job_openings_rate")
            if jolts is not None:
                macro_items.append(
                    (
                        "Job Openings Rate",
                        f"{jolts}%" if isinstance(jolts, (int, float)) else str(jolts),
                    )
                )
            # S50: Additional FRED indicators when primary 3 slots not full
            if len(macro_items) < 3:
                _job_openings = macro_data.get("job_openings")
                if _job_openings is not None and isinstance(
                    _job_openings, (int, float)
                ):
                    macro_items.append(("Job Openings (000s)", f"{_job_openings:,.0f}"))
            if len(macro_items) < 3:
                _ahe = macro_data.get("avg_hourly_earnings")
                if _ahe is not None and isinstance(_ahe, (int, float)):
                    # S3: FRED (US Federal Reserve) macro data is US-only --
                    # mark it explicitly on the figure rather than a bare $.
                    _ahe_val = f"${_ahe:.2f}"
                    if _get_active_currency() != "USD":
                        _ahe_val = _mark_usd(_ahe_val)
                    macro_items.append(("Avg Hourly Earnings (USD)", _ahe_val))
            if len(macro_items) < 3:
                _ffr = macro_data.get("fed_funds_rate")
                if _ffr is not None and isinstance(_ffr, (int, float)):
                    macro_items.append(("Fed Funds Rate", f"{_ffr}%"))

        if not macro_items:
            macro_items = [
                ("Unemployment Rate", "Data not available"),
                ("Labor Force Participation", "Data not available"),
            ]

        macro_card_top = macro_top + Inches(0.4)
        card_w = Inches(1.8)
        card_h = Inches(0.7)
        card_gap = Inches(0.15)

        for mi, (m_label, m_val) in enumerate(macro_items[:3]):
            mx = left_col_left + mi * (card_w + card_gap)
            _add_rounded_rect(slide, mx, macro_card_top, card_w, card_h, WHITE)
            _add_filled_rect(slide, mx, macro_card_top, card_w, Inches(0.04), TEAL)
            _add_textbox(
                slide,
                mx + Inches(0.1),
                macro_card_top + Inches(0.08),
                card_w - Inches(0.2),
                Inches(0.35),
                text=m_val,
                font_size=16,
                bold=True,
                color=NAVY,
                alignment=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE,
            )
            _add_textbox(
                slide,
                mx + Inches(0.1),
                macro_card_top + Inches(0.42),
                card_w - Inches(0.2),
                Inches(0.22),
                text=m_label,
                font_size=7,
                color=MUTED_TEXT,
                alignment=PP_ALIGN.CENTER,
            )

        # ---- RIGHT COLUMN: Salary Intelligence ----
        right_col_left = Inches(7.0)
        right_col_w = Inches(5.8)

        _add_textbox(
            slide,
            right_col_left,
            section_top,
            right_col_w,
            Inches(0.35),
            text="SALARY INTELLIGENCE",
            font_size=11,
            bold=True,
            color=NAVY,
        )
        _add_filled_rect(
            slide,
            right_col_left,
            section_top + Inches(0.33),
            Inches(2.0),
            Inches(0.03),
            TEAL,
        )

        salary_card_top = section_top + Inches(0.5)
        sal_card_h = Inches(1.1)
        sal_card_gap = Inches(0.15)
        sal_card_w = right_col_w

        sal_count = 0
        for role_name, role_sal in list(salary_intel.items())[:4]:
            if not isinstance(role_sal, dict):
                continue
            median = role_sal.get("median") or 0
            if not median or median <= 0:
                continue

            sy = salary_card_top + sal_count * (sal_card_h + sal_card_gap)
            _add_rounded_rect(slide, right_col_left, sy, sal_card_w, sal_card_h, WHITE)
            _add_filled_rect(slide, right_col_left, sy, Inches(0.06), sal_card_h, BLUE)

            # Role name
            _add_textbox(
                slide,
                right_col_left + Inches(0.2),
                sy + Inches(0.06),
                sal_card_w - Inches(0.3),
                Inches(0.25),
                text=str(role_name)[:35],
                font_size=10,
                bold=True,
                color=DARK_TEXT,
            )

            # Salary bar visualization
            sal_min = role_sal.get("min", role_sal.get("p25", median * 0.7))
            sal_max = role_sal.get("max", role_sal.get("p75", median * 1.3))
            sources = role_sal.get("source_count") or 0
            confidence = role_sal.get("confidence") or ""

            bar_left = right_col_left + Inches(0.2)
            bar_top_y = sy + Inches(0.38)
            bar_w = sal_card_w - Inches(0.4)
            bar_h_sal = Inches(0.22)

            # Background bar
            _add_rounded_rect(slide, bar_left, bar_top_y, bar_w, bar_h_sal, LIGHT_BLUE)

            # Median marker (proportional position)
            if sal_max > sal_min and sal_max > 0:
                median_pct = min(
                    1.0, max(0.0, (median - sal_min) / (sal_max - sal_min))
                )
                marker_x = bar_left + bar_w * median_pct - Inches(0.05)
                _add_filled_rect(
                    slide,
                    marker_x,
                    bar_top_y - Inches(0.02),
                    Inches(0.1),
                    bar_h_sal + Inches(0.04),
                    BLUE,
                )

            # Labels
            label_y = sy + Inches(0.65)
            _add_textbox(
                slide,
                bar_left,
                label_y,
                Inches(1.5),
                Inches(0.2),
                text=f"Min: {_format_salary(sal_min)}" if sal_min > 0 else "",
                font_size=7,
                color=MUTED_TEXT,
            )
            _add_textbox(
                slide,
                bar_left + Inches(1.8),
                label_y,
                Inches(1.8),
                Inches(0.2),
                text=f"Median: {_format_salary(median)}",
                font_size=8,
                bold=True,
                color=NAVY,
                alignment=PP_ALIGN.CENTER,
            )
            _add_textbox(
                slide,
                bar_left + Inches(3.5),
                label_y,
                Inches(1.5),
                Inches(0.2),
                text=f"Max: {_format_salary(sal_max)}" if sal_max > 0 else "",
                font_size=7,
                color=MUTED_TEXT,
                alignment=PP_ALIGN.RIGHT,
            )

            # Source/confidence badge
            badge_text = ""
            if sources and sources > 0:
                badge_text = f"{sources} sources"
            if confidence:
                badge_text += f" | {confidence}" if badge_text else str(confidence)
            if badge_text:
                _add_textbox(
                    slide,
                    right_col_left + sal_card_w - Inches(1.8),
                    sy + Inches(0.06),
                    Inches(1.6),
                    Inches(0.2),
                    text=badge_text,
                    font_size=7,
                    italic=True,
                    color=MUTED_TEXT,
                    alignment=PP_ALIGN.RIGHT,
                )

            sal_count += 1

        if sal_count == 0:
            _add_textbox(
                slide,
                right_col_left,
                salary_card_top,
                sal_card_w,
                Inches(0.4),
                text="Salary data not available for selected roles",
                font_size=10,
                italic=True,
                color=MUTED_TEXT,
            )

        # ---- Workforce Trend Highlights (bottom right) ----
        wf_top = (
            section_top
            + Inches(0.5)
            + max(sal_count, 1) * (sal_card_h + sal_card_gap)
            + Inches(0.15)
        )
        _add_textbox(
            slide,
            right_col_left,
            wf_top,
            right_col_w,
            Inches(0.3),
            text="WORKFORCE TREND HIGHLIGHTS",
            font_size=10,
            bold=True,
            color=NAVY,
        )
        _add_filled_rect(
            slide,
            right_col_left,
            wf_top + Inches(0.28),
            Inches(2.2),
            Inches(0.03),
            TEAL,
        )

        wf_bullet_top = wf_top + Inches(0.4)
        wf_bullets = []

        # Gen-Z insights
        gen_z = workforce.get("gen_z_insights", {})
        if isinstance(gen_z, dict):
            wf_share = gen_z.get("workforce_share")
            if wf_share:
                wf_bullets.append(f"Gen-Z now represents {wf_share} of the workforce")
            platforms = gen_z.get("job_search_platforms", {})
            if isinstance(platforms, dict) and platforms:
                top_platform = next(iter(platforms.items()), (None, None))
                if top_platform[0]:
                    wf_bullets.append(
                        f"Top Gen-Z job search: {top_platform[0]} ({top_platform[1]})"
                    )

        # Employer branding
        eb = workforce.get("employer_branding", {})
        if isinstance(eb, dict):
            roi = eb.get("roi_data", {})
            if isinstance(roi, dict) and roi:
                cost_reduction = roi.get("cost_per_hire_reduction")
                if cost_reduction:
                    wf_bullets.append(
                        f"Strong employer brand reduces cost-per-hire by {cost_reduction}"
                    )

        # Research highlights
        research = workforce.get("relevant_research") or []
        if isinstance(research, list):
            for rr in research[:2]:
                if isinstance(rr, dict):
                    title = rr.get("title") or ""
                    publisher = rr.get("publisher") or ""
                    if title:
                        wf_bullets.append(
                            f"Research: {title[:50]}{'...' if len(title) > 50 else ''} ({publisher})"
                        )

        if not wf_bullets:
            wf_bullets = ["Workforce trend data not available for this industry"]

        box_wf, tf_wf = _add_textbox(
            slide, right_col_left, wf_bullet_top, right_col_w, Inches(1.2)
        )
        tf_wf.paragraphs[0].space_before = Pt(0)
        tf_wf.paragraphs[0].space_after = Pt(0)

        for bi, bullet in enumerate(wf_bullets[:4]):
            if bi == 0:
                p = tf_wf.paragraphs[0]
            else:
                p = tf_wf.add_paragraph()
            p.space_before = Pt(2)
            p.space_after = Pt(4)
            rb = p.add_run()
            rb.text = "\u25b8  "
            _set_font(rb, size=9, color=TEAL)
            rt = p.add_run()
            rt.text = str(bullet)
            _set_font(rt, size=8, color=DARK_TEXT)

        # Source line
        _add_textbox(
            slide,
            Inches(0.55),
            Inches(6.7),
            Inches(12.2),
            Inches(0.2),
            text="Sources: Government labor statistics, industry benchmarks, market intelligence",
            font_size=7,
            italic=True,
            color=MUTED_TEXT,
        )

        # Footer
        _add_footer(slide, today)

    except Exception as exc:
        # If slide generation fails, log but don't crash the whole deck
        import logging

        logging.getLogger(__name__).warning("Market analysis slide failed: %s", exc)


# ===================================================================
# SLIDE - Location Analysis (NEW)
# ===================================================================


def _build_slide_location_analysis(prs: Presentation, data: Dict):
    """Build Location Analysis slide using location_profiles data.

    Uses:
    - location_profiles: population, cost of living, regional intelligence,
      top job boards, hiring regulations, cultural norms
    """
    try:
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        client = data.get("client_name", "Client")
        locations = data.get("locations") or []
        today = datetime.date.today().strftime("%B %d, %Y")

        synthesized = data.get("_synthesized", {})
        if not isinstance(synthesized, dict):
            synthesized = {}
        loc_profiles = synthesized.get("location_profiles", {})
        if not isinstance(loc_profiles, dict):
            loc_profiles = {}

        # Fallback: if no synthesized location profiles, build from research.COUNTRY_DATA
        if not loc_profiles and research is not None and locations:
            for loc_str in locations[:4]:
                if not isinstance(loc_str, str):
                    continue
                # Try to detect a country name from the location string
                country_name = research._detect_country(loc_str)
                if country_name and country_name in research.COUNTRY_DATA:
                    cd = research.COUNTRY_DATA[country_name]
                    # Build a profile matching the expected location_profiles schema
                    pop_str = cd.get("population") or ""
                    try:
                        pop_val = int(
                            re.sub(
                                r"[^\d]",
                                "",
                                str(pop_str)
                                .replace("M", "000000")
                                .replace("B", "000000000"),
                            )
                        )
                    except (ValueError, TypeError):
                        pop_val = 0
                    loc_profiles[country_name] = {
                        "population": pop_val,
                        "median_household_income": cd.get("median_salary") or 0,
                        "cost_of_living_index": cd.get("coli") or 0,
                        "currency": cd.get("currency") or "",
                        "timezone": "",
                        "top_job_boards": cd.get("top_boards") or "",
                        "unemployment_rate": cd.get("unemployment") or "",
                        "top_industries": cd.get("top_industries") or "",
                    }
                elif not country_name:
                    # US location -- add a basic "United States" card if not already added
                    if "United States" not in loc_profiles:
                        us_data = research.COUNTRY_DATA.get("United States", {})
                        loc_profiles[loc_str] = {
                            "population": 333000000,
                            "median_household_income": us_data.get(
                                "median_salary", 65000
                            ),
                            "cost_of_living_index": us_data.get("coli", 100),
                            "currency": "USD",
                            "timezone": "",
                            "top_job_boards": us_data.get("top_boards") or "",
                            "unemployment_rate": us_data.get("unemployment") or "",
                            "top_industries": us_data.get("top_industries") or "",
                        }

        # Off-white background
        _add_filled_rect(
            slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, OFF_WHITE
        )

        # Top band
        _add_top_band(slide, "LOCATION ANALYSIS", today)

        n_locs = len(locations)
        action_text = (
            f"Regional market intelligence across {n_locs} target location{'s' if n_locs != 1 else ''} "
            f"for {client}'s recruitment strategy"
        )
        _add_textbox(
            slide,
            Inches(0.55),
            Inches(0.92),
            Inches(12.2),
            Inches(0.5),
            text=action_text,
            font_size=15,
            bold=True,
            color=NAVY,
        )

        # ---- Location Cards ----
        card_top = Inches(1.6)
        total_w = Inches(12.2)
        max_cards = min(len(loc_profiles), 4)  # Show up to 4 locations

        if max_cards == 0:
            # No location data -- show placeholder
            _add_textbox(
                slide,
                Inches(0.55),
                card_top,
                total_w,
                Inches(1.0),
                text="Location profile data not yet available. API enrichment in progress.",
                font_size=14,
                italic=True,
                color=MUTED_TEXT,
                alignment=PP_ALIGN.CENTER,
            )
            _add_footer(slide, today)
            return

        card_gap = Inches(0.2)
        card_w = (
            (total_w - card_gap * (max_cards - 1)) / max_cards
            if max_cards > 1
            else total_w
        )
        card_h = Inches(4.8)

        for li, (loc_name, loc_data) in enumerate(
            list(loc_profiles.items())[:max_cards]
        ):
            if not isinstance(loc_data, dict):
                continue

            cx = Inches(0.55) + li * (card_w + card_gap)

            # Card background
            _add_rounded_rect(slide, cx, card_top, card_w, card_h, WHITE)

            # Location name header
            _add_filled_rect(slide, cx, card_top, card_w, Inches(0.45), NAVY)
            _add_textbox(
                slide,
                cx + Inches(0.15),
                card_top + Inches(0.05),
                card_w - Inches(0.3),
                Inches(0.35),
                text=str(loc_name)[:30],
                font_size=11,
                bold=True,
                color=WHITE,
                anchor=MSO_ANCHOR.MIDDLE,
            )

            content_top = card_top + Inches(0.55)
            content_left = cx + Inches(0.15)
            content_w = card_w - Inches(0.3)

            # Demographics section
            items = []
            pop = loc_data.get("population") or 0
            if pop and pop > 0:
                items.append(("Population", f"{pop:,}"))
            income = loc_data.get("median_household_income") or 0
            if income and income > 0:
                # S3: median_household_income traces only to US Census/DataUSA/
                # METRO_DATA (research.py METRO_DATA is US-metro-only) -- a
                # fixed US-benchmark constant. Mark the figure itself.
                _income_val = f"${income:,}"
                if _get_active_currency() != "USD":
                    _income_val = _mark_usd(_income_val)
                items.append(("Median Income (USD)", _income_val))
            col_index = loc_data.get("cost_of_living_index") or 0
            if col_index and col_index > 0:
                items.append(("Cost of Living", f"{col_index:.0f}/100"))
            talent_density = loc_data.get("talent_density") or 0
            if talent_density and talent_density > 0:
                items.append(("Talent Density", f"{talent_density:.1%}"))
            unemployment = loc_data.get("unemployment_rate") or ""
            if unemployment:
                items.append(("Unemployment", str(unemployment)))
            currency = loc_data.get("currency") or ""
            if currency:
                items.append(("Currency", str(currency)))
            timezone = loc_data.get("timezone") or ""
            if timezone:
                items.append(("Timezone", str(timezone)[:18]))
            top_boards = loc_data.get("top_job_boards") or ""
            if top_boards:
                items.append(("Top Boards", str(top_boards)[:60]))

            box_loc, tf_loc = _add_textbox(
                slide, content_left, content_top, content_w, Inches(1.8)
            )
            tf_loc.paragraphs[0].space_before = Pt(0)
            tf_loc.paragraphs[0].space_after = Pt(0)

            first = True
            for label, value in items[:6]:
                if first:
                    p = tf_loc.paragraphs[0]
                    first = False
                else:
                    p = tf_loc.add_paragraph()
                p.space_before = Pt(1)
                p.space_after = Pt(3)
                rl = p.add_run()
                rl.text = f"{label}: "
                _set_font(rl, size=8, bold=True, color=DARK_TEXT)
                rv = p.add_run()
                rv.text = str(value)
                _set_font(rv, size=8, color=MUTED_TEXT)

            # Regional Intelligence section
            reg_intel = loc_data.get("regional_intelligence", {})
            if isinstance(reg_intel, dict) and reg_intel:
                ri_top = content_top + Inches(1.9)
                _add_filled_rect(
                    slide, content_left, ri_top, content_w, Inches(0.03), TEAL
                )

                _add_textbox(
                    slide,
                    content_left,
                    ri_top + Inches(0.08),
                    content_w,
                    Inches(0.2),
                    text="REGIONAL INTEL",
                    font_size=7,
                    bold=True,
                    color=TEAL,
                )

                ri_items = []

                # Top job boards
                boards = reg_intel.get("top_job_boards") or []
                if isinstance(boards, list) and boards:
                    board_names = [
                        b.get("name", str(b)) if isinstance(b, dict) else str(b)
                        for b in boards[:3]
                    ]
                    ri_items.append(("Top Boards", ", ".join(board_names)))

                # Hiring regulations
                regs = reg_intel.get("hiring_regulations", {})
                if isinstance(regs, dict) and regs:
                    notice_period = regs.get("notice_period") or ""
                    if notice_period:
                        ri_items.append(("Notice Period", str(notice_period)))
                    probation = regs.get("probation_period") or ""
                    if probation:
                        ri_items.append(("Probation", str(probation)))

                # Cultural norms
                norms = reg_intel.get("cultural_norms", {})
                if isinstance(norms, dict) and norms:
                    lang = norms.get("primary_language", norms.get("language") or "")
                    if lang:
                        ri_items.append(("Language", str(lang)))
                    comm = norms.get("communication_style") or ""
                    if comm:
                        ri_items.append(("Comm. Style", str(comm)[:20]))

                # CPA benchmark
                cpa_bench = reg_intel.get("cpa_benchmark", {})
                if isinstance(cpa_bench, dict):
                    cpa_range = cpa_bench.get("range", cpa_bench.get("typical") or "")
                    if cpa_range:
                        ri_items.append(("CPA Range", str(cpa_range)))

                box_ri, tf_ri = _add_textbox(
                    slide, content_left, ri_top + Inches(0.32), content_w, Inches(1.8)
                )
                tf_ri.paragraphs[0].space_before = Pt(0)
                tf_ri.paragraphs[0].space_after = Pt(0)

                ri_first = True
                for rl_label, rl_val in ri_items[:5]:
                    if ri_first:
                        p = tf_ri.paragraphs[0]
                        ri_first = False
                    else:
                        p = tf_ri.add_paragraph()
                    p.space_before = Pt(1)
                    p.space_after = Pt(3)
                    rl_run = p.add_run()
                    rl_run.text = f"\u25b8 {rl_label}: "
                    _set_font(rl_run, size=7, bold=True, color=TEAL)
                    rv_run = p.add_run()
                    rv_run.text = str(rl_val)
                    _set_font(rv_run, size=7, color=DARK_TEXT)

        # Source line
        _add_textbox(
            slide,
            Inches(0.55),
            Inches(6.7),
            Inches(12.2),
            Inches(0.2),
            text="Sources: Government census data, geographic intelligence, international economic data",
            font_size=7,
            italic=True,
            color=MUTED_TEXT,
        )

        # Footer
        _add_footer(slide, today)

    except Exception as exc:
        import logging

        logging.getLogger(__name__).warning("Location analysis slide failed: %s", exc)


# ===================================================================
# SLIDE - Competitive Landscape (NEW)
# ===================================================================

_COMPETITOR_TAG_RE = re.compile(r"^\((?:National|Regional|Local)\)\s*")


def _strip_competitor_tag(name: str) -> str:
    """Strip a leading "(National)"/"(Regional)"/"(Local)" scope tag off a
    competitor name before it's interpolated into a full sentence.

    visual:atria#2: the tag is legitimate as a label on the competitor
    CARD TITLE (it tells the reader this employer's data isn't
    city-specific), but leaking it verbatim into generated prose produces
    a broken-looking sentence: "Counter: (National) Amazon is actively
    competing for..." The card title keeps the tag; only prose
    interpolation strips it.
    """
    return _COMPETITOR_TAG_RE.sub("", str(name or "")).strip() or str(name or "")


# copy:both#1/#2, visual:atria#2: a handful of mega hourly employers (Amazon,
# Walmart, ...) show up in curated/brief competitor lists across every
# vertical -- they DO compete for the same wage band/labor pool, but they
# are not a same-vertical operator (Amazon isn't a senior-living competitor
# for Atria the way Brookdale is) and don't compete for the client's actual
# customers/residents. Per-industry carve-outs cover the sectors where the
# SAME name genuinely is a same-vertical player.
#
# fix/plan-quality-8 (prod Atria bundle): the national logistics/parcel
# carriers (UPS, FedEx, DHL, USPS) belong in this same bucket -- they were
# previously missing, so a senior-living/healthcare client's competitor
# list tagged UPS/FedEx as an "Industry competitor" (same-vertical operator)
# when they only compete for the same wage band/labor pool, not the same
# residents/customers. logistics_supply_chain is carved out below so they
# still read as "Industry competitor" for an actual logistics client.
_CROSS_INDUSTRY_MEGA_EMPLOYERS = frozenset(
    {
        "amazon",
        "walmart",
        "target",
        "costco",
        "home depot",
        "kroger",
        "starbucks",
        "mcdonald's",
        "mcdonalds",
        "ups",
        "fedex",
        "dhl",
        "usps",
    }
)
_INDUSTRY_MEGA_EMPLOYER_EXCEPTIONS: Dict[str, frozenset] = {
    "retail_consumer": frozenset(
        {"amazon", "walmart", "target", "costco", "home depot", "kroger"}
    ),
    "general_entry_level": frozenset(
        {"amazon", "walmart", "target", "costco", "home depot", "kroger"}
    ),
    "hospitality_travel": frozenset({"starbucks", "mcdonald's", "mcdonalds"}),
    "logistics_supply_chain": frozenset(
        {"ups", "fedex", "dhl", "usps", "amazon", "walmart"}
    ),
}


def _classify_competitor_vertical(comp_name: str, industry: str) -> str:
    """Return ``"industry"`` (same-vertical operator -- competes for
    candidates AND the client's own customers/residents) or
    ``"talent_market"`` (a cross-industry hourly employer that only
    competes for the same wage band/labor pool) for a named competitor.

    copy:both#1/#2, visual:atria#2: distinguishes real same-vertical
    competitors (FedEx/UPS for a logistics client, Brookdale for senior
    living) from generic mega-employer filler (Amazon showing up in a
    senior-living competitor list) that reads as unvetted when presented
    in the identical card format as a verified direct competitor.
    """
    name_l = _strip_competitor_tag(str(comp_name or "")).strip().lower()
    if name_l in _CROSS_INDUSTRY_MEGA_EMPLOYERS:
        exceptions = _INDUSTRY_MEGA_EMPLOYER_EXCEPTIONS.get(
            str(industry or "").strip().lower(), frozenset()
        )
        if name_l not in exceptions:
            return "talent_market"
    return "industry"


_COMPETITOR_WHY_TEMPLATES_INDUSTRY = (
    "{name} actively recruits {pool}, competing directly for this plan's "
    "pipeline -- and for the same customers.",
    "{name} is a {type_phrase}same-vertical employer drawing from "
    "the same {pool} this plan targets.",
    "{name}'s hiring activity for similar roles puts direct pressure on "
    "{pool}, while also competing for the same customer base.",
    "{name} is visibly hiring in the same market as {pool}, a same-vertical "
    "source of candidate AND customer overlap.",
)
_COMPETITOR_WHY_TEMPLATES_TALENT_MARKET = (
    "{name} isn't a same-vertical operator, but it draws from the same "
    "labor pool and wage band as {pool}.",
    "{name} is a {type_phrase}large hourly employer competing for "
    "the same wage band as {pool}, not the same customers.",
    "{name}'s hourly pay and hiring volume put wage-band pressure on "
    "{pool}, without competing for this industry's customers.",
    "{name} pulls from the same local labor pool as {pool} -- a wage-band "
    "competitor, not a same-vertical one.",
)


def _compose_competitor_why(
    comp_name: str, comp_data: Dict, ctx: Dict, ordinal: int
) -> str:
    """One "why this competitor matters" sentence, varied by ``ordinal`` and
    composed from competitor type + role + city -- never the single
    "Competing employer in <city>" sentence repeated byte-identically for
    every card (copy:both#3 / visual:atria#2: the previous "Why:" text was
    the SAME line for every competitor pulled from one city, and
    industry-agnostic for brief-supplied competitors with no description).

    copy:both#1/#2: also varies by ``ctx["vertical_type"]`` -- a
    "talent_market" competitor (e.g. Amazon next to two real senior-living
    operators) gets labor-pool/wage-band framing instead of the same
    "competing for candidates" language used for a verified same-vertical
    competitor.
    """
    role = str(ctx.get("role") or "").strip()
    city = str(ctx.get("city") or "").strip()
    if role and city:
        pool = f"{role} candidates in {city}"
    elif role:
        pool = f"{role} candidates"
    elif city:
        pool = f"candidates in {city}"
    else:
        pool = "this plan's candidate pool"

    ctype = str((comp_data or {}).get("competitor_type") or "").strip().lower()
    type_phrase = (
        f"{ctype} " if ctype and ctype not in ("national", "regional", "local") else ""
    )
    # "an employer" when no type qualifier precedes it, "a <type> employer"
    # otherwise (grammar-safe for the common case; type qualifiers are
    # curated lowercase words like "senior_living", not free text).
    article = "an" if not type_phrase else "a"

    name = _strip_competitor_tag(comp_name) or "This competitor"
    vertical_type = str(ctx.get("vertical_type") or "industry").strip().lower()
    templates = (
        _COMPETITOR_WHY_TEMPLATES_TALENT_MARKET
        if vertical_type == "talent_market"
        else _COMPETITOR_WHY_TEMPLATES_INDUSTRY
    )
    idx = (
        int(ordinal) % len(templates)
        if isinstance(ordinal, (int, float)) and not isinstance(ordinal, bool)
        else 0
    )
    return templates[idx].format(
        name=name, pool=pool, type_phrase=type_phrase, article=article
    )


# fix/plan-quality-8 (prod Atria bundle): the competitor card column starts
# at section_top(1.6in) + 0.5in = 2.1in and each card (with the taller
# counter-strategy prose) is 1.3in + a 0.12in gap = 1.42in tall. The footer
# rule sits at y=7.12in (_add_footer). 3 cards bottom out at
# 2.1 + 3*1.3 + 2*0.12 = 6.24in -- comfortably above the footer. A 4th card
# bottoms out at 7.66in, past both the footer band and the 7.5in slide
# itself. Cap the rendered (and collected) competitor count here so a 4th
# card never gets drawn partially off the slide.
_MAX_COMPETITOR_CARDS = 3


def _build_slide_competitive_landscape(prs: Presentation, data: Dict):
    """Build the Competitive Landscape slide.

    Uses:
    - competitive_intelligence: company profile, competitor data,
      industry hiring trends, market positioning
    """
    try:
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        client = data.get("client_name", "Client")
        industry_label = data.get("industry_label") or ""
        roles = data.get("roles") or []
        locations = data.get("locations") or []
        _cl_first_role = str(roles[0]) if roles else ""
        _cl_first_loc = (
            str(locations[0])
            if locations and isinstance(locations[0], (str, int, float))
            else ""
        )
        today = datetime.date.today().strftime("%B %d, %Y")

        synthesized = data.get("_synthesized", {})
        if not isinstance(synthesized, dict):
            synthesized = {}
        comp_intel = synthesized.get("competitive_intelligence", {})
        if not isinstance(comp_intel, dict):
            comp_intel = {}

        # Fallback: if no competitive intelligence from synthesis, try knowledge base
        if not comp_intel:
            kb = data.get("_knowledge_base", {})
            if isinstance(kb, dict) and kb:
                industry_key = data.get("industry", "general_entry_level")
                # Try recruitment_benchmarks section for industry-level data
                rb = kb.get("recruitment_benchmarks", {})
                if isinstance(rb, dict):
                    ind_bench = rb.get("industry_benchmarks", {}).get(industry_key, {})
                    if not isinstance(ind_bench, dict):
                        # Try alternative key formats (e.g., "technology_engineering" vs "tech_engineering")
                        for kb_key in rb.get("industry_benchmarks", {}):
                            if industry_key.split("_")[0] in kb_key:
                                ind_bench = rb["industry_benchmarks"][kb_key]
                                break
                    if ind_bench:
                        # Build a minimal comp_intel from KB benchmarks
                        hiring_trends_fb = {}
                        if ind_bench.get("time_to_fill"):
                            hiring_trends_fb["avg_time_to_fill"] = ind_bench[
                                "time_to_fill"
                            ]
                        if ind_bench.get("offer_acceptance_rate"):
                            hiring_trends_fb["offer_acceptance_rate"] = ind_bench[
                                "offer_acceptance_rate"
                            ]
                        if ind_bench.get("quality_of_hire"):
                            hiring_trends_fb["quality_metrics"] = ind_bench[
                                "quality_of_hire"
                            ]
                        if ind_bench.get("source_of_hire"):
                            hiring_trends_fb["top_sources"] = ind_bench[
                                "source_of_hire"
                            ]
                        if hiring_trends_fb:
                            comp_intel["hiring_trends"] = hiring_trends_fb
                            comp_intel["company_profile"] = {"name": client}

        # Off-white background
        _add_filled_rect(
            slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, OFF_WHITE
        )

        # Top band
        _add_top_band(slide, "COMPETITIVE LANDSCAPE", today)

        action_text = (
            f"Market positioning and competitor intelligence for "
            f"{client}'s talent acquisition strategy in {industry_label}"
        )
        _add_textbox(
            slide,
            Inches(0.55),
            Inches(0.92),
            Inches(12.2),
            Inches(0.5),
            text=action_text,
            font_size=15,
            bold=True,
            color=NAVY,
        )

        # ---- LEFT: Company Profile ----
        section_top = Inches(1.6)
        left_w = Inches(5.5)

        _add_textbox(
            slide,
            Inches(0.55),
            section_top,
            left_w,
            Inches(0.35),
            text="COMPANY PROFILE",
            font_size=11,
            bold=True,
            color=NAVY,
        )
        _add_filled_rect(
            slide,
            Inches(0.55),
            section_top + Inches(0.33),
            Inches(1.8),
            Inches(0.03),
            TEAL,
        )

        company = comp_intel.get("company_profile", {})
        if not isinstance(company, dict):
            company = {}

        profile_items = [
            ("Company", company.get("name", client)),
        ]
        desc = company.get("description") or ""
        if desc:
            # Word-boundary truncation (never mid-word), ellipsis only when
            # the text was actually cut.
            profile_items.append(("Description", _trunc_word(str(desc), 100)))
        domain = company.get("domain") or ""
        if domain:
            profile_items.append(("Domain", str(domain)))
        is_public = company.get("is_public", False)
        if is_public:
            ticker = company.get("sec_ticker") or ""
            profile_items.append(
                ("Public Company", f"Ticker: {ticker}" if ticker else "Yes")
            )
            filings = company.get("recent_filings_count") or 0
            if filings and filings > 0:
                profile_items.append(("SEC Filings", f"{filings} recent filings"))
        sic_desc = company.get("sec_sic_description") or ""
        if sic_desc:
            profile_items.append(("SIC Industry", str(sic_desc)[:50]))
        tags = company.get("clearbit_tags") or []
        if isinstance(tags, list) and tags:
            profile_items.append(("Tags", ", ".join(str(t) for t in tags[:4])))
        profile_items = profile_items[:7]

        # Card shrinks to fit the number of POPULATED fields (1-7 rows) --
        # never a fixed-height card with mostly blank space below a
        # 1-2-field profile.
        _profile_row_h_in = 0.27
        profile_card_h = Inches(
            max(0.62, min(2.2, 0.3 + len(profile_items) * _profile_row_h_in))
        )

        profile_top = section_top + Inches(0.5)
        _add_rounded_rect(slide, Inches(0.55), profile_top, left_w, profile_card_h, WHITE)
        _add_filled_rect(
            slide, Inches(0.55), profile_top, Inches(0.06), profile_card_h, BLUE
        )

        box_p, tf_p = _add_textbox(
            slide,
            Inches(0.8),
            profile_top + Inches(0.15),
            left_w - Inches(0.4),
            profile_card_h - Inches(0.2),
        )
        tf_p.paragraphs[0].space_before = Pt(0)
        tf_p.paragraphs[0].space_after = Pt(0)

        first = True
        for label, value in profile_items:
            if first:
                p = tf_p.paragraphs[0]
                first = False
            else:
                p = tf_p.add_paragraph()
            p.space_before = Pt(2)
            p.space_after = Pt(4)
            rl = p.add_run()
            rl.text = f"{label}:  "
            _set_font(rl, size=9, bold=True, color=DARK_TEXT)
            rv = p.add_run()
            rv.text = str(value)
            _set_font(rv, size=9, color=MUTED_TEXT)

        # ---- Industry Hiring Trends (below company profile) ----
        trends_top = profile_top + profile_card_h + Inches(0.2)
        _add_textbox(
            slide,
            Inches(0.55),
            trends_top,
            left_w,
            Inches(0.3),
            text="INDUSTRY HIRING TRENDS",
            font_size=10,
            bold=True,
            color=NAVY,
        )
        _add_filled_rect(
            slide,
            Inches(0.55),
            trends_top + Inches(0.28),
            Inches(2.0),
            Inches(0.03),
            TEAL,
        )

        hiring_trends = comp_intel.get("hiring_trends", {})
        if not isinstance(hiring_trends, dict):
            hiring_trends = {}

        trend_items = []
        emp_count = hiring_trends.get("employment_count")
        if emp_count and isinstance(emp_count, (int, float)) and emp_count > 0:
            trend_items.append(f"Industry employment: {int(emp_count):,}")
        emp_growth = hiring_trends.get("employment_growth_rate")
        if emp_growth is not None:
            trend_items.append(f"Growth rate: {emp_growth}")
        avg_wage = hiring_trends.get("average_weekly_wage")
        if avg_wage and isinstance(avg_wage, (int, float)) and avg_wage > 0:
            # S3: sourced from BLS QCEW (US Bureau of Labor Statistics) --
            # a fixed US-benchmark constant. Mark the figure itself.
            _wage_str = f"Avg weekly wage: ${avg_wage:,.0f}"
            if _get_active_currency() != "USD":
                _wage_str = _mark_usd(_wage_str)
            trend_items.append(_wage_str)
        establishments = hiring_trends.get("establishments")
        if (
            establishments
            and isinstance(establishments, (int, float))
            and establishments > 0
        ):
            trend_items.append(f"Establishments: {int(establishments):,}")

        # KB-derived trends
        kb_insights = hiring_trends.get("kb_insights", {})
        if isinstance(kb_insights, dict):
            outlook = kb_insights.get("outlook") or ""
            if outlook:
                trend_items.append(f"Outlook: {outlook}")
            demand_drivers = kb_insights.get("demand_drivers") or []
            if isinstance(demand_drivers, list) and demand_drivers:
                trend_items.append(
                    f"Drivers: {', '.join(str(d) for d in demand_drivers[:3])}"
                )

        # KB benchmark fallback data (from _knowledge_base)
        ttf = hiring_trends.get("avg_time_to_fill")
        if ttf:
            if isinstance(ttf, dict):
                # Format dict: prefer average_days, then join key-value pairs (skip notes)
                avg_d = ttf.get("average_days") or ttf.get("average") or ""
                parts = []
                if avg_d:
                    parts.append(
                        f"{avg_d} days"
                        if isinstance(avg_d, (int, float))
                        else str(avg_d)
                    )
                for k, v in ttf.items():
                    if k in ("average_days", "average", "notes") or not v:
                        continue
                    # copy:both#5-family: plain .title() clobbers acronyms
                    # embedded in KB keys (e.g. "cdl_drivers" -> "Cdl
                    # Drivers" instead of "CDL Drivers").
                    label = _fmt.smart_title(k)
                    parts.append(f"{label}: {v}")
                ttf_str = " | ".join(parts[:4]) if parts else str(ttf)
            else:
                ttf_str = str(ttf)
            trend_items.append(f"Avg Time-to-Fill: {ttf_str}")
        oar = hiring_trends.get("offer_acceptance_rate")
        if oar:
            oar_str = (
                str(oar)
                if not isinstance(oar, dict)
                else " | ".join(
                    f"{_fmt.smart_title(k)}: {v}"
                    for k, v in oar.items()
                    if k != "notes" and v
                )
            )
            trend_items.append(f"Offer Acceptance: {oar_str}")
        top_src = hiring_trends.get("top_sources")
        if isinstance(top_src, dict) and top_src:
            # S92 fix: this dict's keys come straight from the KB JSON
            # (e.g. "job_boards", "career_site", "employee_referrals") --
            # title-case them like the ttf/oar dict-flattening two blocks
            # above already does, instead of leaking the raw snake_case key.
            src_items = [
                f"{_fmt.smart_title(k)}: {v}"
                for k, v in list(top_src.items())[:3]
            ]
            trend_items.append(f"Top Sources: {', '.join(src_items)}")

        # Fewer than 2 real bullets: fill with genuine, sourced substance
        # (curated role-requirement callouts + a city-data-backed geography
        # rationale) instead of a bare "not available" placeholder.
        if len(trend_items) < 2:
            _cl_gold = data.get("_gold_standard") or {}
            _cl_city_data = (
                _cl_gold.get("city_level_data") or {} if isinstance(_cl_gold, dict) else {}
            )
            trend_items.extend(
                _insight.role_requirements_callout(str(data.get("industry") or ""), roles)
            )
            if _cl_first_loc:
                trend_items.append(
                    _insight.geography_rationale(
                        _cl_first_loc, _cl_city_data.get(_cl_first_loc.split(",")[0].strip())
                    )
                )

        if not trend_items:
            trend_items = ["Industry trend data not available"]

        box_t, tf_t = _add_textbox(
            slide, Inches(0.55), trends_top + Inches(0.4), left_w, Inches(1.5)
        )
        tf_t.paragraphs[0].space_before = Pt(0)
        tf_t.paragraphs[0].space_after = Pt(0)

        for ti, item in enumerate(trend_items[:5]):
            if ti == 0:
                p = tf_t.paragraphs[0]
            else:
                p = tf_t.add_paragraph()
            p.space_before = Pt(1)
            p.space_after = Pt(4)
            rb = p.add_run()
            rb.text = "\u25b8  "
            _set_font(rb, size=9, color=TEAL)
            rt = p.add_run()
            rt.text = str(item)
            _set_font(rt, size=9, color=DARK_TEXT)

        # ---- RIGHT: Competitor Cards ----
        right_left = Inches(6.5)
        right_w = Inches(6.3)

        _add_textbox(
            slide,
            right_left,
            section_top,
            right_w,
            Inches(0.35),
            text="COMPETITOR LANDSCAPE",
            font_size=11,
            bold=True,
            color=NAVY,
        )
        _add_filled_rect(
            slide,
            right_left,
            section_top + Inches(0.33),
            Inches(2.2),
            Inches(0.03),
            TEAL,
        )

        # visual:atria#2 / copy:both#3: brief-supplied competitors are the
        # planner's own explicit, vetted competitive set for THIS client and
        # role -- they must win over any generic/national-default source
        # (a real brief naming "Brookdale/Sunrise/Amazon" was previously
        # being overridden by a generic "Amazon/Walmart/UPS" national-
        # employer list because the brief was checked LAST in the cascade).
        # Only fall through to synthesized/gold-standard data when the
        # brief didn't supply any competitors of its own.
        competitors: Dict[str, Any] = {}
        _direct_comps = data.get("competitors") or []
        if isinstance(_direct_comps, list) and _direct_comps:
            for _dc in _direct_comps[:_MAX_COMPETITOR_CARDS]:
                _dc_name = (
                    str(_dc).strip()
                    if isinstance(_dc, str)
                    else (
                        str(_dc.get("name", "")).strip()
                        if isinstance(_dc, dict)
                        else ""
                    )
                )
                if _dc_name:
                    competitors[_dc_name] = {
                        "domain": (
                            _dc.get("domain", "") if isinstance(_dc, dict) else ""
                        ),
                        "description": (
                            _dc.get("description", "") if isinstance(_dc, dict) else ""
                        ),
                        "competitor_type": (
                            _dc.get("competitor_type", "")
                            if isinstance(_dc, dict)
                            else ""
                        ),
                    }
        elif isinstance(_direct_comps, dict) and _direct_comps:
            competitors = dict(list(_direct_comps.items())[:_MAX_COMPETITOR_CARDS])

        if not competitors:
            _ci_competitors = comp_intel.get("competitors", {})
            if isinstance(_ci_competitors, dict):
                competitors = dict(_ci_competitors)

        # Fallback: pull from gold standard competitor_mapping if still empty.
        # consistency:manpower#5 / visual:atria#2(d): competitor_mapping is
        # keyed by CITY and covers every market the plan targets (the same
        # multi-market table excel_v2's Quality Intelligence "Competitive
        # Landscape" section already renders in full) -- round-robin ONE top
        # employer per city instead of draining the first city's top-3
        # employers and never reaching the plan's other markets.
        if not competitors:
            gold = data.get("_gold_standard") or {}
            comp_map = gold.get("competitor_mapping") or {}
            if isinstance(comp_map, dict) and comp_map:
                _city_employer_lists = [
                    (city_name, city_info.get("top_employers") or [])
                    for city_name, city_info in comp_map.items()
                    if isinstance(city_name, str)
                    and not city_name.startswith("_")
                    and isinstance(city_info, dict)
                ]
                seen: set = set()
                _rr_idx = 0
                while len(competitors) < _MAX_COMPETITOR_CARDS and _city_employer_lists:
                    _added_this_round = False
                    for city_name, employers in _city_employer_lists:
                        if len(competitors) >= _MAX_COMPETITOR_CARDS:
                            break
                        if _rr_idx < len(employers):
                            emp_name = str(employers[_rr_idx]).strip()
                            if emp_name and emp_name not in seen:
                                seen.add(emp_name)
                                competitors[emp_name] = {
                                    "domain": "",
                                    "description": "",
                                    "_market": city_name,
                                }
                                _added_this_round = True
                    if not _added_this_round:
                        break
                    _rr_idx += 1

        if not competitors:
            gold = data.get("_gold_standard") or {}
            _gs_comp_list = gold.get("competitors_list") or []
            if isinstance(_gs_comp_list, list) and _gs_comp_list:
                for _gc in _gs_comp_list[:_MAX_COMPETITOR_CARDS]:
                    _gc_name = (
                        str(_gc).strip()
                        if isinstance(_gc, str)
                        else (
                            str(_gc.get("name", "")).strip()
                            if isinstance(_gc, dict)
                            else ""
                        )
                    )
                    if _gc_name:
                        competitors[_gc_name] = {
                            "domain": "",
                            "description": "",
                        }

        comp_card_top = section_top + Inches(0.5)
        comp_card_h = Inches(0.8)
        comp_card_gap = Inches(0.12)

        if not competitors:
            _add_textbox(
                slide,
                right_left,
                comp_card_top,
                right_w,
                Inches(0.4),
                text="No competitor data available. Add competitors to your request.",
                font_size=10,
                italic=True,
                color=MUTED_TEXT,
            )
        else:
            # Taller cards for counter-strategy -- compose_counter_strategy's
            # prose runs longer than the old fixed boilerplate line, so the
            # card needs more room than the original 1.0in fit.
            comp_card_h = Inches(1.3)
            comp_card_gap = Inches(0.12)
            for ci, (comp_name, comp_data) in enumerate(
                list(competitors.items())[:_MAX_COMPETITOR_CARDS]
            ):
                if not isinstance(comp_data, dict):
                    continue
                cy = comp_card_top + ci * (comp_card_h + comp_card_gap)

                _add_rounded_rect(slide, right_left, cy, right_w, comp_card_h, WHITE)

                # Competitor name with color accent
                accent_colors = [BLUE, TEAL, NAVY, GREEN, AMBER]
                accent = accent_colors[ci % len(accent_colors)]
                _add_filled_rect(
                    slide, right_left, cy, Inches(0.06), comp_card_h, accent
                )

                _add_textbox(
                    slide,
                    right_left + Inches(0.2),
                    cy + Inches(0.05),
                    Inches(3.0),
                    Inches(0.25),
                    text=str(comp_name),
                    font_size=10,
                    bold=True,
                    color=DARK_TEXT,
                )

                # copy:both#1/#2, visual:atria#2: tag each card as a real
                # same-vertical "Industry competitor" or a generic
                # "Talent-market competitor" (e.g. Amazon in a senior-living
                # list) so a non-industry name doesn't read as a vetted
                # direct competitor.
                _vertical_type = _classify_competitor_vertical(
                    comp_name, str(data.get("industry") or "")
                )
                _tag_text, _tag_color = (
                    ("Talent-market competitor", AMBER)
                    if _vertical_type == "talent_market"
                    else ("Industry competitor", MUTED_TEXT)
                )
                _add_textbox(
                    slide,
                    right_left + Inches(3.3),
                    cy + Inches(0.05),
                    right_w - Inches(3.5),
                    Inches(0.25),
                    text=_tag_text,
                    font_size=7,
                    italic=True,
                    color=_tag_color,
                    alignment=PP_ALIGN.RIGHT,
                    anchor=MSO_ANCHOR.MIDDLE,
                )

                # Counter-strategy context -- built once, shared by the Why
                # and Counter lines so both draw from the same per-market
                # signal (comp_data["_market"], set by the round-robin
                # multi-market fallback) instead of always defaulting to
                # the plan's first location.
                _comp_market = str(comp_data.get("_market") or "").strip()
                _counter_ctx = {
                    "role": _cl_first_role,
                    "city": _comp_market or _cl_first_loc,
                    "industry": industry_label,
                    "intensity": str(comp_data.get("hiring_intensity") or ""),
                    "competitor_type": str(comp_data.get("competitor_type") or ""),
                    "vertical_type": _vertical_type,
                    "ordinal": ci,
                }

                # Why they matter -- copy:both#3 / visual:atria#2(c): a
                # composed, per-competitor sentence (varies by ordinal +
                # competitor type/role/city) instead of one repeated
                # boilerplate line. A genuinely specific description (real
                # API-sourced content, not one of the known generic
                # placeholders) is still preferred when present.
                comp_domain = comp_data.get("domain") or ""
                comp_desc = str(comp_data.get("description") or "").strip()
                _generic_desc = comp_desc.lower() in (
                    "",
                    "competitor",
                    "industry competitor",
                ) or comp_desc.lower().startswith("competing employer in")
                if comp_desc and not _generic_desc:
                    why_text = _trunc_word(comp_desc, 80)
                elif comp_domain:
                    why_text = f"Competes for the same talent pool via {comp_domain}"
                else:
                    why_text = _compose_competitor_why(
                        comp_name, comp_data, _counter_ctx, ci
                    )

                _add_textbox(
                    slide,
                    right_left + Inches(0.2),
                    cy + Inches(0.3),
                    right_w - Inches(0.4),
                    Inches(0.25),
                    text=f"Why: {why_text}",
                    font_size=8,
                    color=MUTED_TEXT,
                )

                # Counter-strategy -- per-competitor prose (varies
                # deterministically by competitor + role), never the same
                # boilerplate sentence for every card. visual:atria#2(b):
                # strip any "(National)" scope tag before interpolating the
                # name into prose -- it stays on the card title above.
                counter = _trunc_word(
                    "Counter: "
                    + _insight.compose_counter_strategy(
                        _strip_competitor_tag(comp_name), _counter_ctx
                    ),
                    190,
                )
                _add_textbox(
                    slide,
                    right_left + Inches(0.2),
                    cy + Inches(0.58),
                    right_w - Inches(0.4),
                    Inches(0.65),
                    text=counter,
                    font_size=8,
                    bold=False,
                    color=GREEN,
                )

        # Market positioning insight
        positioning = comp_intel.get("market_positioning", {})
        if isinstance(positioning, dict) and positioning:
            pos_text = positioning.get("summary", positioning.get("insight") or "")
            # Only draw the callout band when there's actual text to show --
            # an empty highlight band (colored rect with nothing in it) used
            # to render whenever `positioning` was a non-empty dict with no
            # summary/insight field.
            if pos_text:
                pos_top = Inches(5.8)
                _add_rounded_rect(
                    slide, Inches(0.55), pos_top, Inches(12.2), Inches(0.7), PALE_TEAL
                )
                _add_filled_rect(
                    slide, Inches(0.55), pos_top, Inches(0.06), Inches(0.7), TEAL
                )
                _add_textbox(
                    slide,
                    Inches(0.8),
                    pos_top + Inches(0.1),
                    Inches(11.7),
                    Inches(0.5),
                    text=_trunc_word(str(pos_text), 200),
                    font_size=9,
                    color=DARK_TEXT,
                )

        # Source line
        _add_textbox(
            slide,
            Inches(0.55),
            Inches(6.7),
            Inches(12.2),
            Inches(0.2),
            text="Sources: Public company filings, industry employment data, curated knowledge base",
            font_size=7,
            italic=True,
            color=MUTED_TEXT,
        )

        # Footer
        _add_footer(slide, today)

    except Exception as exc:
        import logging

        logging.getLogger(__name__).warning(
            "Competitive landscape slide failed: %s", exc
        )


# ===================================================================
# SLIDE - Workforce Trends (NEW)
# ===================================================================


def _build_slide_geopolitical_risk(prs: Presentation, data: Dict):
    """Build the Geopolitical Risk Assessment slide.

    Uses _synthesized.geopolitical_context from api_enrichment.
    Shows risk score badge, per-location event cards, recommendations.
    Only shown when risk_level is not 'low'.
    """
    try:
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        client = data.get("client_name", "Client")
        today = datetime.date.today().strftime("%B %d, %Y")
        synthesized = data.get("_synthesized", {})
        if not isinstance(synthesized, dict):
            synthesized = {}
        geo = synthesized.get("geopolitical_context", {})
        if not isinstance(geo, dict):
            geo = {}

        risk_level = geo.get("risk_level", "low")
        overall_score = geo.get("overall_risk_score", 1.0)
        summary_text = geo.get(
            "summary", "No significant geopolitical risks identified."
        )
        recommendations = geo.get("recommendations") or []
        loc_data = geo.get("locations", {})

        # Background
        _add_filled_rect(
            slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, OFF_WHITE
        )

        # Top band - use red for high/critical, amber for moderate
        band_color = NAVY
        if risk_level in ("high", "critical"):
            band_color = RED_ACCENT
        elif risk_level == "moderate":
            band_color = AMBER
        _add_top_band(
            slide, "GEOPOLITICAL RISK ASSESSMENT", today, band_color=band_color
        )

        # Action text
        _add_textbox(
            slide,
            Inches(0.55),
            Inches(0.92),
            Inches(12.2),
            Inches(0.5),
            text=f"Market risk factors impacting {client}'s recruitment campaigns",
            font_size=15,
            bold=True,
            color=NAVY,
        )

        # Risk score badge (top right area)
        badge_color = (
            GREEN
            if overall_score <= 3
            else (AMBER if overall_score <= 6 else RED_ACCENT)
        )
        badge_bg = (
            LIGHT_GREEN
            if overall_score <= 3
            else (LIGHT_AMBER if overall_score <= 6 else RGBColor(0xFD, 0xE2, 0xE2))
        )
        _add_rounded_rect(
            slide, Inches(10.5), Inches(0.85), Inches(2.2), Inches(0.6), badge_bg
        )
        _add_textbox(
            slide,
            Inches(10.5),
            Inches(0.88),
            Inches(2.2),
            Inches(0.55),
            text=f"Risk: {risk_level.upper()} ({overall_score:.1f}/10)",
            font_size=12,
            bold=True,
            color=badge_color,
            alignment=PP_ALIGN.CENTER,
        )

        # Summary section -- size text to fit longer narratives
        summary_display = _trunc_word(summary_text, 500)
        summary_font = 11 if len(summary_display) <= 200 else 9
        summary_box_h = Inches(0.9) if len(summary_display) <= 200 else Inches(1.1)
        _add_rounded_rect(
            slide, Inches(0.55), Inches(1.6), Inches(12.2), summary_box_h, WHITE
        )
        _add_textbox(
            slide,
            Inches(0.75),
            Inches(1.65),
            Inches(11.8),
            summary_box_h - Inches(0.1),
            text=summary_display,
            font_size=summary_font,
            color=DARK_TEXT,
        )

        # Per-location cards -- adjust top for variable summary height
        card_top = Inches(1.6) + summary_box_h + Inches(0.2)
        max_locations = min(len(loc_data), 4)
        if max_locations > 0:
            card_w = Inches((12.0 / max_locations) - 0.15)
            card_gap = Inches(0.15)

            for i, (loc_name, loc_info) in enumerate(list(loc_data.items())[:4]):
                if not isinstance(loc_info, dict):
                    continue
                card_left = Inches(0.55) + i * (card_w + card_gap)
                loc_score = loc_info.get("risk_score", 1.0)
                loc_events = loc_info.get("events") or []
                loc_adj = loc_info.get("budget_adjustment_factor", 1.0)

                # Card background
                card_bg = (
                    LIGHT_GREEN
                    if loc_score <= 3
                    else (LIGHT_AMBER if loc_score <= 6 else RGBColor(0xFD, 0xE2, 0xE2))
                )
                _add_rounded_rect(
                    slide, card_left, card_top, card_w, Inches(2.8), WHITE
                )
                # Top accent bar
                accent = (
                    GREEN
                    if loc_score <= 3
                    else (AMBER if loc_score <= 6 else RED_ACCENT)
                )
                _add_filled_rect(
                    slide, card_left, card_top, card_w, Inches(0.05), accent
                )

                # Location name + score
                _add_textbox(
                    slide,
                    card_left + Inches(0.1),
                    card_top + Inches(0.1),
                    card_w - Inches(0.2),
                    Inches(0.3),
                    text=str(loc_name)[:30],
                    font_size=11,
                    bold=True,
                    color=NAVY,
                )
                _add_textbox(
                    slide,
                    card_left + Inches(0.1),
                    card_top + Inches(0.38),
                    card_w - Inches(0.2),
                    Inches(0.25),
                    text=f"Risk: {loc_score:.1f}/10  |  Budget adj: {loc_adj:.2f}x",
                    font_size=9,
                    color=accent,
                    bold=True,
                )

                # Events list
                event_y = card_top + Inches(0.7)
                for ev in loc_events[:3]:
                    if not isinstance(ev, dict):
                        continue
                    ev_text = ev.get("event") or ""[:80]
                    severity = ev.get("severity", "low")
                    sev_icon = {
                        "low": " ",
                        "moderate": " ",
                        "high": " ",
                        "critical": " ",
                    }.get(severity, " ")
                    _add_textbox(
                        slide,
                        card_left + Inches(0.1),
                        event_y,
                        card_w - Inches(0.2),
                        Inches(0.45),
                        text=f"{sev_icon}{ev_text}",
                        font_size=8,
                        color=DARK_TEXT,
                    )
                    event_y += Inches(0.42)

        # Recommendations section
        rec_top = Inches(5.7)
        if recommendations:
            _add_textbox(
                slide,
                Inches(0.55),
                rec_top,
                Inches(12.2),
                Inches(0.3),
                text="RECOMMENDATIONS",
                font_size=10,
                bold=True,
                color=BLUE,
            )

            rec_text = " | ".join(_trunc_word(r, 200) for r in recommendations[:4])
            _add_textbox(
                slide,
                Inches(0.55),
                rec_top + Inches(0.3),
                Inches(12.2),
                Inches(0.9),
                text=rec_text,
                font_size=9,
                color=DARK_TEXT,
            )

        # Source attribution -- sanitize internal provider names
        _llm_keywords = {
            "claude",
            "haiku",
            "sonnet",
            "opus",
            "gpt",
            "gemini",
            "groq",
            "llama",
            "mistral",
            "anthropic",
            "openai",
        }
        raw_source = str(geo.get("source") or "")
        # Strip provider names and confidence suffixes from source string
        if any(kw in raw_source.lower() for kw in _llm_keywords) or not raw_source:
            source_display = "AI Analysis"
        else:
            source_display = raw_source
        # Format confidence as proper percentage
        raw_conf = geo.get("confidence") or 0
        try:
            conf_val = float(raw_conf)
        except (TypeError, ValueError):
            conf_val = 0.0
        if 0 < conf_val <= 1.0:
            conf_val = conf_val * 100
        conf_str = f"{conf_val:.0f}%" if conf_val > 0 else "—"
        _add_textbox(
            slide,
            Inches(0.55),
            Inches(6.8),
            Inches(12.2),
            Inches(0.3),
            text=f"Source: {source_display} | Confidence: {conf_str}",
            font_size=7,
            italic=True,
            color=MUTED_TEXT,
        )

    except Exception as exc:
        logger.warning("Failed to build geopolitical risk slide: %s", exc)


def _build_slide_workforce_trends(prs: Presentation, data: Dict):
    """Build the Workforce Trends slide.

    Uses:
    - workforce_insights: Gen-Z preferences, employer branding,
      white paper citations, remote work trends, supply partner trends
    """
    try:
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        client = data.get("client_name", "Client")
        industry_label = data.get("industry_label") or ""
        today = datetime.date.today().strftime("%B %d, %Y")

        synthesized = data.get("_synthesized", {})
        if not isinstance(synthesized, dict):
            synthesized = {}
        workforce = synthesized.get("workforce_insights", {})
        if not isinstance(workforce, dict):
            workforce = {}

        # Off-white background
        _add_filled_rect(
            slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, OFF_WHITE
        )

        # Top band
        _add_top_band(slide, "WORKFORCE TRENDS & INSIGHTS", today)

        action_text = (
            f"Emerging workforce trends shaping {client}'s talent acquisition "
            f"strategy in {industry_label}"
        )
        _add_textbox(
            slide,
            Inches(0.55),
            Inches(0.92),
            Inches(12.2),
            Inches(0.5),
            text=action_text,
            font_size=15,
            bold=True,
            color=NAVY,
        )

        # ---- THREE-COLUMN LAYOUT ----
        section_top = Inches(1.6)
        col_w = Inches(3.95)
        col_gap = Inches(0.2)
        col1_left = Inches(0.55)
        col2_left = col1_left + col_w + col_gap
        col3_left = col2_left + col_w + col_gap
        col_h = Inches(4.5)

        # ---- COLUMN 1: Gen-Z Insights ----
        _add_rounded_rect(slide, col1_left, section_top, col_w, col_h, WHITE)
        _add_filled_rect(slide, col1_left, section_top, col_w, Inches(0.05), BLUE)

        _add_textbox(
            slide,
            col1_left + Inches(0.15),
            section_top + Inches(0.12),
            col_w - Inches(0.3),
            Inches(0.3),
            text="GEN-Z WORKFORCE TRENDS",
            font_size=10,
            bold=True,
            color=BLUE,
        )

        gen_z = workforce.get("gen_z_insights", {})
        if not isinstance(gen_z, dict):
            gen_z = {}

        gz_items = []

        wf_share = gen_z.get("workforce_share")
        if wf_share:
            gz_items.append(("Workforce Share", str(wf_share)))

        # Platform preferences
        platforms = gen_z.get("job_search_platforms", {})
        if isinstance(platforms, dict) and platforms:
            for pname, pval in list(platforms.items())[:3]:
                gz_items.append((str(pname), str(pval)))

        # Mobile vs desktop
        mobile = gen_z.get("mobile_vs_desktop", {})
        if isinstance(mobile, dict):
            mobile_pct = mobile.get("mobile", mobile.get("mobile_first") or "")
            if mobile_pct:
                gz_items.append(("Mobile Usage", str(mobile_pct)))

        # Social media habits
        social = gen_z.get("social_media_habits", {})
        if isinstance(social, dict) and social:
            for sname, sval in list(social.items())[:2]:
                gz_items.append((str(sname).title(), str(sval)))

        # Workplace expectations
        expectations = gen_z.get("workplace_expectations", {})
        if isinstance(expectations, dict):
            flex = expectations.get("flexibility", {})
            if isinstance(flex, dict):
                remote_pref = flex.get(
                    "remote_preference", flex.get("flexible_work") or ""
                )
                if remote_pref:
                    gz_items.append(("Flexibility", str(remote_pref)))
            dei = expectations.get("dei", {})
            if isinstance(dei, dict):
                dei_imp = dei.get("importance", dei.get("priority") or "")
                if dei_imp:
                    gz_items.append(("DEI Expectations", str(dei_imp)))
            mh = expectations.get("mental_health", {})
            if isinstance(mh, dict):
                mh_priority = mh.get("priority", mh.get("importance") or "")
                if mh_priority:
                    gz_items.append(("Mental Health", str(mh_priority)))

        # Tenure
        tenure = gen_z.get("tenure", {})
        if isinstance(tenure, dict):
            avg_tenure = tenure.get("average", tenure.get("median") or "")
            if avg_tenure:
                gz_items.append(("Avg Tenure", str(avg_tenure)))

        if not gz_items:
            gz_items = [("Status", "Gen-Z data not available")]

        box_gz, tf_gz = _add_textbox(
            slide,
            col1_left + Inches(0.15),
            section_top + Inches(0.5),
            col_w - Inches(0.3),
            col_h - Inches(0.6),
        )
        tf_gz.paragraphs[0].space_before = Pt(0)
        tf_gz.paragraphs[0].space_after = Pt(0)

        for gi, (g_label, g_val) in enumerate(gz_items[:10]):
            if gi == 0:
                p = tf_gz.paragraphs[0]
            else:
                p = tf_gz.add_paragraph()
            p.space_before = Pt(2)
            p.space_after = Pt(4)
            rl = p.add_run()
            rl.text = f"{g_label}:  "
            _set_font(rl, size=8, bold=True, color=DARK_TEXT)
            rv = p.add_run()
            rv.text = str(g_val)[:50]
            _set_font(rv, size=8, color=MUTED_TEXT)

        # ---- COLUMN 2: Employer Branding ----
        _add_rounded_rect(slide, col2_left, section_top, col_w, col_h, WHITE)
        _add_filled_rect(slide, col2_left, section_top, col_w, Inches(0.05), TEAL)

        _add_textbox(
            slide,
            col2_left + Inches(0.15),
            section_top + Inches(0.12),
            col_w - Inches(0.3),
            Inches(0.3),
            text="EMPLOYER BRANDING",
            font_size=10,
            bold=True,
            color=TEAL,
        )

        eb = workforce.get("employer_branding", {})
        if not isinstance(eb, dict):
            eb = {}

        eb_items = []

        # ROI data
        roi = eb.get("roi_data", {})
        if isinstance(roi, dict):
            for rk, rv_val in list(roi.items())[:5]:
                label = str(rk).replace("_", " ").title()
                eb_items.append((label, str(rv_val)))

        # Best practices
        bp = eb.get("best_practices", {})
        if isinstance(bp, dict):
            for bk, bv in list(bp.items())[:3]:
                label = str(bk).replace("_", " ").title()
                if isinstance(bv, list):
                    eb_items.append((label, ", ".join(str(v) for v in bv[:3])))
                elif isinstance(bv, dict):
                    eb_items.append((label, str(next(iter(bv.values()), ""))))
                else:
                    eb_items.append((label, str(bv)[:50]))

        # Channel effectiveness
        ch_eff = eb.get("channel_effectiveness", {})
        if isinstance(ch_eff, dict) and ch_eff:
            for ck, cv in list(ch_eff.items())[:3]:
                label = str(ck).replace("_", " ").title()
                if isinstance(cv, dict):
                    eb_items.append(
                        (f"Channel: {label}", str(next(iter(cv.values()), "")))
                    )
                else:
                    eb_items.append((f"Channel: {label}", str(cv)[:40]))

        if not eb_items:
            eb_items = [("Status", "Employer branding data not available")]

        box_eb, tf_eb = _add_textbox(
            slide,
            col2_left + Inches(0.15),
            section_top + Inches(0.5),
            col_w - Inches(0.3),
            col_h - Inches(0.6),
        )
        tf_eb.paragraphs[0].space_before = Pt(0)
        tf_eb.paragraphs[0].space_after = Pt(0)

        for ei, (e_label, e_val) in enumerate(eb_items[:10]):
            if ei == 0:
                p = tf_eb.paragraphs[0]
            else:
                p = tf_eb.add_paragraph()
            p.space_before = Pt(2)
            p.space_after = Pt(4)
            rl = p.add_run()
            rl.text = f"{e_label}:  "
            _set_font(rl, size=8, bold=True, color=DARK_TEXT)
            rv = p.add_run()
            rv.text = str(e_val)[:50]
            _set_font(rv, size=8, color=MUTED_TEXT)

        # ---- COLUMN 3: Research & Supply Trends ----
        _add_rounded_rect(slide, col3_left, section_top, col_w, col_h, WHITE)
        _add_filled_rect(slide, col3_left, section_top, col_w, Inches(0.05), NAVY)

        _add_textbox(
            slide,
            col3_left + Inches(0.15),
            section_top + Inches(0.12),
            col_w - Inches(0.3),
            Inches(0.3),
            text="RESEARCH & INDUSTRY DATA",
            font_size=10,
            bold=True,
            color=NAVY,
        )

        # White paper citations
        research = workforce.get("relevant_research") or []
        r_items = []
        if isinstance(research, list):
            for rr in research[:4]:
                if isinstance(rr, dict):
                    title = rr.get("title") or ""
                    publisher = rr.get("publisher") or ""
                    year = rr.get("year") or ""
                    findings = rr.get("top_findings") or []
                    if title:
                        r_items.append(
                            {
                                "title": str(title)[:60],
                                "publisher": str(publisher),
                                "year": str(year) if year else "",
                                "finding": (
                                    str(findings[0])[:60]
                                    if isinstance(findings, list) and findings
                                    else ""
                                ),
                            }
                        )

        # Supply partner trends
        sp = workforce.get("supply_partner_trends", {})
        # Job type trends
        jt = workforce.get("job_type_trends", {})

        box_r, tf_r = _add_textbox(
            slide,
            col3_left + Inches(0.15),
            section_top + Inches(0.5),
            col_w - Inches(0.3),
            col_h - Inches(0.6),
        )
        tf_r.paragraphs[0].space_before = Pt(0)
        tf_r.paragraphs[0].space_after = Pt(0)

        first_r = True
        for ri_item in r_items:
            if first_r:
                p = tf_r.paragraphs[0]
                first_r = False
            else:
                p = tf_r.add_paragraph()
            p.space_before = Pt(3)
            p.space_after = Pt(2)
            rt = p.add_run()
            rt.text = ri_item["title"]
            _set_font(rt, size=8, bold=True, color=DARK_TEXT)

            p2 = tf_r.add_paragraph()
            p2.space_before = Pt(0)
            p2.space_after = Pt(2)
            rs = p2.add_run()
            source_str = ri_item["publisher"]
            if ri_item["year"]:
                source_str += f" ({ri_item['year']})"
            rs.text = source_str
            _set_font(rs, size=7, italic=True, color=MUTED_TEXT)

            if ri_item["finding"]:
                p3 = tf_r.add_paragraph()
                p3.space_before = Pt(0)
                p3.space_after = Pt(4)
                rf = p3.add_run()
                rf.text = f"\u25b8 {ri_item['finding']}"
                _set_font(rf, size=7, color=TEAL)

        # Supply partner trends section
        if isinstance(sp, dict) and sp:
            p_sp = tf_r.add_paragraph()
            p_sp.space_before = Pt(6)
            p_sp.space_after = Pt(2)
            r_sp = p_sp.add_run()
            r_sp.text = "Supply Partner Trends"
            _set_font(r_sp, size=8, bold=True, color=NAVY)

            for sk, sv in list(sp.items())[:3]:
                p_s = tf_r.add_paragraph()
                p_s.space_before = Pt(1)
                p_s.space_after = Pt(3)
                rs1 = p_s.add_run()
                rs1.text = f"\u25b8 {str(sk).replace('_', ' ').title()}: "
                _set_font(rs1, size=7, bold=True, color=TEAL)
                rs2 = p_s.add_run()
                if isinstance(sv, dict):
                    rs2.text = str(next(iter(sv.values()), ""))[:40]
                else:
                    rs2.text = str(sv)[:40]
                _set_font(rs2, size=7, color=DARK_TEXT)

        # Job type trends section
        if isinstance(jt, dict) and jt:
            p_jt = tf_r.add_paragraph()
            p_jt.space_before = Pt(6)
            p_jt.space_after = Pt(2)
            r_jt = p_jt.add_run()
            r_jt.text = "Job Type Trends"
            _set_font(r_jt, size=8, bold=True, color=NAVY)

            for jk, jv in list(jt.items())[:3]:
                p_j = tf_r.add_paragraph()
                p_j.space_before = Pt(1)
                p_j.space_after = Pt(3)
                rj1 = p_j.add_run()
                rj1.text = f"\u25b8 {str(jk).replace('_', ' ').title()}: "
                _set_font(rj1, size=7, bold=True, color=TEAL)
                rj2 = p_j.add_run()
                if isinstance(jv, dict):
                    rj2.text = str(next(iter(jv.values()), ""))[:40]
                else:
                    rj2.text = str(jv)[:40]
                _set_font(rj2, size=7, color=DARK_TEXT)

        if not r_items and not sp and not jt:
            p = tf_r.paragraphs[0]
            r = p.add_run()
            r.text = "Research and trend data not available for this industry"
            _set_font(r, size=9, italic=True, color=MUTED_TEXT)

        # Source line
        _add_textbox(
            slide,
            Inches(0.55),
            Inches(6.7),
            Inches(12.2),
            Inches(0.2),
            text="Sources: Recruitment Industry White Papers, Workforce Trends Intelligence, Employer Branding Research",
            font_size=7,
            italic=True,
            color=MUTED_TEXT,
        )

        # Footer
        _add_footer(slide, today)

    except Exception as exc:
        import logging

        logging.getLogger(__name__).warning("Workforce trends slide failed: %s", exc)


# ===================================================================
# SLIDE: Creative Testing Plan (A/B Test Recommendations)
# ===================================================================


# -- Channel-specific creative testing data --
_CHANNEL_CREATIVE_DATA: Dict[str, Dict[str, Any]] = {
    "indeed": {
        "variants": [
            "Benefit-focused: Lead with salary range and benefits package",
            "Growth-focused: Emphasize career development and training",
            "Culture-focused: Highlight team dynamics and work environment",
        ],
        "bidding": "Start with CPC bidding; shift to CPA once 50+ conversions",
        "ctr_range": "2.5% - 4.2%",
    },
    "linkedin": {
        "variants": [
            "Professional tone: Industry-specific language and requirements",
            "Employer brand: Company mission and values-led messaging",
            "Opportunity-focused: Unique projects and impact statements",
        ],
        "bidding": "Maximum delivery bidding for awareness; target cost for applications",
        "ctr_range": "0.4% - 0.9%",
    },
    "google": {
        "variants": [
            "Direct CTA: 'Apply Now' with role title and location",
            "Question-based: 'Looking for [Role]? Join [Company]'",
            "Benefit-led: '[Salary] + Benefits - [Role] at [Company]'",
        ],
        "bidding": "Enhanced CPC with conversion tracking; shift to target CPA at scale",
        "ctr_range": "3.0% - 6.5%",
    },
    "facebook": {
        "variants": [
            "Video testimonial: 30s employee day-in-the-life",
            "Carousel: Multiple roles with distinct creative per card",
            "Single image: Bold text overlay with clear CTA",
        ],
        "bidding": "Lowest cost bidding; use cost cap once baseline CPA established",
        "ctr_range": "0.8% - 2.1%",
    },
    "programmatic": {
        "variants": [
            "Retargeting: Personalized ads for previous job page visitors",
            "Contextual: Industry-relevant placement with role-specific copy",
            "Dynamic: Auto-populate role title and location from feed",
        ],
        "bidding": "Programmatic CPC with daily budget caps; optimize toward CPA",
        "ctr_range": "0.3% - 1.2%",
    },
}

_INDUSTRY_TEST_FRAMEWORKS: Dict[str, List[str]] = {
    "healthcare_medical": [
        "Test credential requirements prominence (above vs. below fold)",
        "Compare sign-on bonus vs. career growth messaging",
        "A/B shift flexibility language (flexible vs. set schedules)",
    ],
    "tech_engineering": [
        "Test remote-first vs. hybrid messaging in job titles",
        "Compare tech stack listing vs. project impact descriptions",
        "A/B salary transparency (range shown vs. competitive compensation)",
    ],
    "retail_consumer": [
        "Test immediate start language vs. standard posting",
        "Compare hourly rate prominence vs. total compensation",
        "A/B employee discount/perks messaging effectiveness",
    ],
    "default": [
        "Test salary transparency vs. competitive compensation language",
        "Compare short-form (3 bullet) vs. detailed job descriptions",
        "A/B apply button placement and CTA wording",
    ],
}


def _build_slide_creative_testing(prs: Presentation, data: Dict) -> None:
    """Build the Creative Testing Plan slide with A/B test recommendations.

    Provides per-channel ad copy variants, testing framework, bidding
    strategies, and expected CTR ranges based on the plan's channel list
    and industry vertical.
    """
    try:
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        client = data.get("client_name") or "Client"
        industry = data.get("industry") or "general_entry_level"
        industry_label = (
            data.get("industry_label") or industry.replace("_", " ").title()
        )
        today = datetime.date.today().strftime("%B %d, %Y")

        # Off-white background
        _add_filled_rect(
            slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, OFF_WHITE
        )

        # Top band
        _add_top_band(slide, "CREATIVE TESTING PLAN", today)

        # Subtitle
        _add_textbox(
            slide,
            Inches(0.55),
            Inches(0.92),
            Inches(12.2),
            Inches(0.5),
            text=f"A/B testing recommendations for {client}'s {industry_label} campaigns",
            font_size=15,
            bold=True,
            color=NAVY,
        )

        # Get channels from plan data
        channels_raw = data.get("channel_categories") or {}
        channel_names: List[str] = []
        if isinstance(channels_raw, dict):
            channel_names = [k for k, v in channels_raw.items() if v]
        elif isinstance(channels_raw, list):
            channel_names = [
                (c.get("name") or str(c)) if isinstance(c, dict) else str(c)
                for c in channels_raw
            ]

        # Map channel names to creative data keys
        matched_channels: List[Tuple[str, Dict]] = []
        for ch_name in channel_names[:5]:
            ch_lower = ch_name.lower().replace("_", " ")
            for key, cdata in _CHANNEL_CREATIVE_DATA.items():
                if key in ch_lower or ch_lower in key:
                    matched_channels.append((_smart_title(ch_name), cdata))
                    break
        # Ensure at least 2 channels
        if len(matched_channels) < 2:
            for key in ["indeed", "linkedin", "google"]:
                if len(matched_channels) >= 3:
                    break
                if not any(key in mc[0].lower() for mc in matched_channels):
                    matched_channels.append((key.title(), _CHANNEL_CREATIVE_DATA[key]))

        # ── LEFT COLUMN: Channel-specific variants (60% width) ──
        col1_left = Inches(0.55)
        col1_w = Inches(7.3)
        col2_left = Inches(8.1)
        col2_w = Inches(4.7)
        section_top = Inches(1.55)

        # Channel variant cards
        _add_rounded_rect(slide, col1_left, section_top, col1_w, Inches(5.0), WHITE)
        _add_filled_rect(slide, col1_left, section_top, col1_w, Inches(0.05), BLUE)

        _add_textbox(
            slide,
            col1_left + Inches(0.15),
            section_top + Inches(0.1),
            col1_w - Inches(0.3),
            Inches(0.3),
            text="AD COPY VARIANTS BY CHANNEL",
            font_size=10,
            bold=True,
            color=BLUE,
        )

        box_ch, tf_ch = _add_textbox(
            slide,
            col1_left + Inches(0.15),
            section_top + Inches(0.45),
            col1_w - Inches(0.3),
            Inches(4.4),
        )
        tf_ch.paragraphs[0].space_before = Pt(0)
        tf_ch.paragraphs[0].space_after = Pt(0)

        first_item = True
        for ch_label, ch_data in matched_channels[:4]:
            if first_item:
                p = tf_ch.paragraphs[0]
                first_item = False
            else:
                p = tf_ch.add_paragraph()
            p.space_before = Pt(6)
            p.space_after = Pt(2)
            rh = p.add_run()
            rh.text = f"{ch_label}  (CTR: {ch_data.get('ctr_range', 'N/A')})"
            _set_font(rh, size=9, bold=True, color=NAVY)

            for vi, variant in enumerate(ch_data.get("variants", [])[:3]):
                pv = tf_ch.add_paragraph()
                pv.space_before = Pt(1)
                pv.space_after = Pt(2)
                rv = pv.add_run()
                rv.text = f"  {chr(65 + vi)}. {variant}"
                _set_font(rv, size=7, color=MUTED_TEXT)

            # Bidding strategy
            pb = tf_ch.add_paragraph()
            pb.space_before = Pt(1)
            pb.space_after = Pt(4)
            rb = pb.add_run()
            rb.text = f"  Bidding: {ch_data.get('bidding', 'CPC recommended')}"
            _set_font(rb, size=7, italic=True, color=TEAL)

        # ── RIGHT COLUMN: Testing Framework ──
        _add_rounded_rect(slide, col2_left, section_top, col2_w, Inches(5.0), WHITE)
        _add_filled_rect(slide, col2_left, section_top, col2_w, Inches(0.05), TEAL)

        _add_textbox(
            slide,
            col2_left + Inches(0.15),
            section_top + Inches(0.1),
            col2_w - Inches(0.3),
            Inches(0.3),
            text="TESTING FRAMEWORK",
            font_size=10,
            bold=True,
            color=TEAL,
        )

        box_fw, tf_fw = _add_textbox(
            slide,
            col2_left + Inches(0.15),
            section_top + Inches(0.45),
            col2_w - Inches(0.3),
            Inches(4.4),
        )
        tf_fw.paragraphs[0].space_before = Pt(0)
        tf_fw.paragraphs[0].space_after = Pt(0)

        # A/B test structure
        framework_items = [
            ("Test Duration", "7-14 days per variant (min. 100 conversions)"),
            ("Sample Split", "50/50 traffic split; 95% confidence threshold"),
            ("Primary Metric", "Cost Per Application (CPA)"),
            ("Secondary Metrics", "CTR, Apply Rate, Quality Score"),
            ("Winner Criteria", "Lower CPA at statistical significance"),
        ]

        first_fw = True
        for fw_label, fw_val in framework_items:
            if first_fw:
                p = tf_fw.paragraphs[0]
                first_fw = False
            else:
                p = tf_fw.add_paragraph()
            p.space_before = Pt(3)
            p.space_after = Pt(3)
            rl = p.add_run()
            rl.text = f"{fw_label}:  "
            _set_font(rl, size=8, bold=True, color=DARK_TEXT)
            rv = p.add_run()
            rv.text = fw_val
            _set_font(rv, size=8, color=MUTED_TEXT)

        # Industry-specific tests
        p_ind_h = tf_fw.add_paragraph()
        p_ind_h.space_before = Pt(10)
        p_ind_h.space_after = Pt(4)
        r_ind_h = p_ind_h.add_run()
        r_ind_h.text = f"INDUSTRY-SPECIFIC TESTS ({industry_label.upper()})"
        _set_font(r_ind_h, size=8, bold=True, color=NAVY)

        industry_tests = _INDUSTRY_TEST_FRAMEWORKS.get(
            industry, _INDUSTRY_TEST_FRAMEWORKS["default"]
        )
        for test_desc in industry_tests[:3]:
            pt = tf_fw.add_paragraph()
            pt.space_before = Pt(2)
            pt.space_after = Pt(3)
            rt = pt.add_run()
            rt.text = f"\u25b8 {test_desc}"
            _set_font(rt, size=7, color=TEAL)

        # Source line
        _add_textbox(
            slide,
            Inches(0.55),
            Inches(6.7),
            Inches(12.2),
            Inches(0.2),
            text="Source: Nova AI Suite Creative Testing Engine | Industry ad platform benchmarks",
            font_size=7,
            italic=True,
            color=MUTED_TEXT,
        )

        # Footer
        _add_footer(slide, today)

    except Exception as exc:
        logger.warning("Failed to build creative testing slide: %s", exc)


# ===================================================================
# GOLD STANDARD SLIDES -- Market Intelligence, Strategy, Timeline
# ===================================================================


def _build_slides_gold_standard(prs: Presentation, data: Dict) -> None:
    """Build 2-3 slides from Gold Standard quality gate outputs.

    Adds the following slides (if data is available):
      - Market Intelligence: city-level data + competitor mapping
      - Strategy & Difficulty: channel strategy + role difficulty framework
      - Activation Timeline: hiring calendar + budget phasing

    Each slide is wrapped in its own try/except so a single failure
    does not block the remaining slides.

    Args:
        prs: The python-pptx Presentation object.
        data: The full enriched data dict containing ``_gold_standard``.
    """
    gold: dict = data.get("_gold_standard") or {}
    if not gold:
        return

    today = datetime.date.today().strftime("%B %d, %Y")

    # ── Slide A: Market Intelligence (city-level + competitors) ──
    city_data: dict = gold.get("city_level_data") or {}
    competitor_map: dict = gold.get("competitor_mapping") or {}
    clearance: Optional[dict] = gold.get("clearance_segmentation")

    if city_data or competitor_map:
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _add_filled_rect(
                slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, OFF_WHITE
            )
            _add_top_band(slide, "MARKET INTELLIGENCE", today)

            _add_textbox(
                slide,
                Inches(0.55),
                Inches(0.92),
                Inches(12.2),
                Inches(0.4),
                text="City-level supply-demand analysis and competitive landscape",
                font_size=14,
                bold=True,
                color=NAVY,
            )

            content_top = Inches(1.45)

            # ---- City-Level Data Table ----
            if city_data:
                _add_textbox(
                    slide,
                    Inches(0.55),
                    content_top,
                    Inches(5.5),
                    Inches(0.3),
                    text="CITY-LEVEL SUPPLY DATA",
                    font_size=10,
                    bold=True,
                    color=NAVY,
                )
                _add_filled_rect(
                    slide,
                    Inches(0.55),
                    content_top + Inches(0.3),
                    Inches(1.5),
                    Inches(0.025),
                    TEAL,
                )

                # Table header
                city_table_top = content_top + Inches(0.45)
                col_labels = ["City", "Salary", "Difficulty", "Supply Tier"]
                col_widths = [Inches(1.5), Inches(1.3), Inches(1.1), Inches(1.3)]
                total_w = sum(w for w in col_widths)

                _add_filled_rect(
                    slide,
                    Inches(0.55),
                    city_table_top,
                    total_w,
                    Inches(0.3),
                    NAVY,
                )
                x_pos = Inches(0.55)
                for label, cw in zip(col_labels, col_widths):
                    _add_textbox(
                        slide,
                        x_pos + Inches(0.05),
                        city_table_top,
                        cw - Inches(0.05),
                        Inches(0.3),
                        text=label,
                        font_size=8,
                        bold=True,
                        color=WHITE,
                        anchor=MSO_ANCHOR.MIDDLE,
                    )
                    x_pos += cw

                # Table rows (limit to 8 cities to fit)
                row_y = city_table_top + Inches(0.3)
                row_h = Inches(0.28)
                cities_shown = list(city_data.items())[:8]
                for idx, (city_name, info) in enumerate(cities_shown):
                    bg = LIGHT_BLUE if idx % 2 == 0 else WARM_WHITE
                    _add_filled_rect(slide, Inches(0.55), row_y, total_w, row_h, bg)
                    vals = [
                        city_name,
                        # S3: derived from the plan's own national_avg_salary
                        # (enrich_city_level_data) -- localize.
                        _fmt_currency(info.get("estimated_salary", 0)),
                        f"{info.get('hiring_difficulty', 0):.1f}/10",
                        str(info.get("supply_tier") or "balanced")
                        .replace("_", " ")
                        .title(),
                    ]
                    x_pos = Inches(0.55)
                    for val, cw in zip(vals, col_widths):
                        _add_textbox(
                            slide,
                            x_pos + Inches(0.05),
                            row_y,
                            cw - Inches(0.05),
                            row_h,
                            text=val,
                            font_size=7,
                            color=DARK_TEXT,
                            anchor=MSO_ANCHOR.MIDDLE,
                        )
                        x_pos += cw
                    row_y += row_h

            # ---- Competitor Mapping (right side) ----
            if competitor_map:
                comp_left = Inches(6.2)
                _add_textbox(
                    slide,
                    comp_left,
                    content_top,
                    Inches(6.5),
                    Inches(0.3),
                    text="COMPETITOR LANDSCAPE",
                    font_size=10,
                    bold=True,
                    color=NAVY,
                )
                _add_filled_rect(
                    slide,
                    comp_left,
                    content_top + Inches(0.3),
                    Inches(1.5),
                    Inches(0.025),
                    TEAL,
                )

                comp_y = content_top + Inches(0.5)
                comp_entries = [
                    (city, info)
                    for city, info in competitor_map.items()
                    if not city.startswith("_")
                ][:6]

                for city_name, info in comp_entries:
                    employers = info.get("top_employers") or []
                    intensity = str(info.get("hiring_intensity") or "moderate").title()

                    _add_textbox(
                        slide,
                        comp_left,
                        comp_y,
                        Inches(6.5),
                        Inches(0.22),
                        text=f"{city_name} ({intensity} intensity)",
                        font_size=8,
                        bold=True,
                        color=NAVY,
                    )
                    _add_textbox(
                        slide,
                        comp_left + Inches(0.1),
                        comp_y + Inches(0.22),
                        Inches(6.4),
                        Inches(0.22),
                        text=", ".join(employers[:5]),
                        font_size=7,
                        color=MUTED_TEXT,
                    )
                    comp_y += Inches(0.5)

            # Clearance badge (if applicable). S4: non-US plans carry
            # clearance_segmentation["us_framework_only"]=True with no
            # primary_clearance (see gold_standard.detect_clearance_requirements)
            # -- never render an "N/A" clearance badge for those; skip it.
            if clearance and not clearance.get("us_framework_only"):
                primary = clearance.get("primary_clearance") or {}
                badge_text = f"Clearance: {primary.get('level', 'N/A')} | +{primary.get('salary_premium_pct', 0)}% premium"
                _add_filled_rect(
                    slide, Inches(0.55), Inches(6.4), Inches(4), Inches(0.35), NAVY
                )
                _add_textbox(
                    slide,
                    Inches(0.65),
                    Inches(6.4),
                    Inches(3.8),
                    Inches(0.35),
                    text=badge_text,
                    font_size=8,
                    bold=True,
                    color=TEAL,
                    anchor=MSO_ANCHOR.MIDDLE,
                )

            _add_footer(slide, today)

        except Exception as exc:
            logger.error(
                "Gold Standard slide (Market Intelligence) failed: %s",
                exc,
                exc_info=True,
            )

    # ── Slide B: Strategy & Difficulty ──
    channel_strategy: dict = gold.get("channel_strategy") or {}
    difficulty_framework: list = gold.get("difficulty_framework") or []

    if channel_strategy or difficulty_framework:
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _add_filled_rect(
                slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, OFF_WHITE
            )
            _add_top_band(slide, "STRATEGY & DIFFICULTY FRAMEWORK", today)

            _add_textbox(
                slide,
                Inches(0.55),
                Inches(0.92),
                Inches(12.2),
                Inches(0.4),
                text="Channel strategy split and role complexity classification",
                font_size=14,
                bold=True,
                color=NAVY,
            )

            content_top = Inches(1.5)

            # ---- Channel Strategy (left half) ----
            if channel_strategy:
                split = channel_strategy.get("recommended_split") or {}
                trad_pct = split.get("traditional_pct", 65)
                nontrad_pct = split.get("non_traditional_pct", 35)

                _add_textbox(
                    slide,
                    Inches(0.55),
                    content_top,
                    Inches(5.5),
                    Inches(0.3),
                    text="CHANNEL MIX STRATEGY",
                    font_size=10,
                    bold=True,
                    color=NAVY,
                )
                _add_filled_rect(
                    slide,
                    Inches(0.55),
                    content_top + Inches(0.3),
                    Inches(1.5),
                    Inches(0.025),
                    TEAL,
                )

                # Split visualization bar
                bar_top = content_top + Inches(0.5)
                bar_h = Inches(0.5)
                bar_w = Inches(5.2)
                trad_w = Inches(5.2 * trad_pct / 100)
                nontrad_w = Inches(5.2 * nontrad_pct / 100)

                _add_filled_rect(slide, Inches(0.55), bar_top, trad_w, bar_h, BLUE)
                _add_textbox(
                    slide,
                    Inches(0.6),
                    bar_top,
                    trad_w - Inches(0.1),
                    bar_h,
                    text=f"Traditional {trad_pct}%",
                    font_size=10,
                    bold=True,
                    color=WHITE,
                    anchor=MSO_ANCHOR.MIDDLE,
                )

                _add_filled_rect(
                    slide, Inches(0.55) + trad_w, bar_top, nontrad_w, bar_h, TEAL
                )
                _add_textbox(
                    slide,
                    Inches(0.6) + trad_w,
                    bar_top,
                    nontrad_w - Inches(0.1),
                    bar_h,
                    text=f"Non-Trad {nontrad_pct}%",
                    font_size=10,
                    bold=True,
                    color=WHITE,
                    anchor=MSO_ANCHOR.MIDDLE,
                )

                # Channel lists below bar
                list_top = bar_top + Inches(0.65)
                trad_channels = channel_strategy.get("traditional_channels") or []
                nontrad_channels = (
                    channel_strategy.get("non_traditional_channels") or []
                )

                for idx, ch in enumerate(trad_channels[:4]):
                    _add_textbox(
                        slide,
                        Inches(0.7),
                        list_top + Inches(idx * 0.22),
                        Inches(2.5),
                        Inches(0.22),
                        text=f"  {ch.get('name', '')}",
                        font_size=7,
                        color=DARK_TEXT,
                    )

                for idx, ch in enumerate(nontrad_channels[:4]):
                    _add_textbox(
                        slide,
                        Inches(3.2),
                        list_top + Inches(idx * 0.22),
                        Inches(2.8),
                        Inches(0.22),
                        text=f"  {ch.get('name', '')}",
                        font_size=7,
                        color=DARK_TEXT,
                    )

            # ---- Difficulty Framework (right half) ----
            if difficulty_framework:
                diff_left = Inches(6.5)
                _add_textbox(
                    slide,
                    diff_left,
                    content_top,
                    Inches(6),
                    Inches(0.3),
                    text="ROLE DIFFICULTY CLASSIFICATION",
                    font_size=10,
                    bold=True,
                    color=NAVY,
                )
                _add_filled_rect(
                    slide,
                    diff_left,
                    content_top + Inches(0.3),
                    Inches(1.5),
                    Inches(0.025),
                    TEAL,
                )

                # Difficulty cards
                card_top = content_top + Inches(0.5)
                card_h = Inches(0.7)
                card_w = Inches(6.2)

                for idx, role_info in enumerate(difficulty_framework[:6]):
                    card_y = card_top + Inches(idx * (0.7 + 0.08))
                    complexity = role_info.get("complexity_score", 0)
                    seniority = str(role_info.get("seniority_level") or "mid").title()

                    # Complexity-based color
                    if complexity >= 7:
                        card_color = LIGHT_AMBER
                    elif complexity >= 4:
                        card_color = LIGHT_BLUE
                    else:
                        card_color = LIGHT_GREEN

                    _add_filled_rect(
                        slide, diff_left, card_y, card_w, card_h, card_color
                    )

                    _add_textbox(
                        slide,
                        diff_left + Inches(0.1),
                        card_y + Inches(0.05),
                        Inches(3),
                        Inches(0.25),
                        text=str(role_info.get("role_title") or ""),
                        font_size=9,
                        bold=True,
                        color=NAVY,
                    )
                    _add_textbox(
                        slide,
                        diff_left + Inches(0.1),
                        card_y + Inches(0.3),
                        Inches(5.8),
                        Inches(0.35),
                        text=(
                            f"{seniority} | Complexity: {complexity}/10 | "
                            f"TTF: {role_info.get('avg_time_to_fill_days', 0)} days | "
                            f"Budget: {role_info.get('budget_weight', 1.0):.1f}x"
                        ),
                        font_size=7,
                        color=MUTED_TEXT,
                    )

            _add_footer(slide, today)

        except Exception as exc:
            logger.error(
                "Gold Standard slide (Strategy & Difficulty) failed: %s",
                exc,
                exc_info=True,
            )

    # ── Slide C: Activation Timeline ──
    activation: dict = gold.get("activation_calendar") or {}
    budget_tiers: dict = gold.get("budget_tiers") or {}

    if activation:
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _add_filled_rect(
                slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, OFF_WHITE
            )
            _add_top_band(slide, "ACTIVATION TIMELINE", today)

            _add_textbox(
                slide,
                Inches(0.55),
                Inches(0.92),
                Inches(12.2),
                Inches(0.4),
                text="Hiring intensity calendar and budget phasing",
                font_size=14,
                bold=True,
                color=NAVY,
            )

            content_top = Inches(1.5)
            timeline: list = activation.get("timeline") or []

            if timeline:
                _add_textbox(
                    slide,
                    Inches(0.55),
                    content_top,
                    Inches(12),
                    Inches(0.3),
                    text="6-MONTH HIRING CALENDAR",
                    font_size=10,
                    bold=True,
                    color=NAVY,
                )
                _add_filled_rect(
                    slide,
                    Inches(0.55),
                    content_top + Inches(0.3),
                    Inches(1.5),
                    Inches(0.025),
                    TEAL,
                )

                # Calendar cards (horizontal layout)
                card_w = Inches(1.95)
                card_h = Inches(3.5)
                cards_left = Inches(0.55)
                cards_top = content_top + Inches(0.5)
                gap = Inches(0.1)

                intensity_colors: dict[str, RGBColor] = {
                    "very_high": RGBColor(0xB7, 0x66, 0x9E),  # Magenta
                    "high": RGBColor(0x5A, 0x54, 0xBE),  # Blue Violet
                    "moderate": RGBColor(0x6B, 0xB5, 0xCE),  # Teal
                    "low": RGBColor(0xEB, 0xE6, 0xE0),  # Warm gray
                }

                for idx, month_info in enumerate(timeline[:6]):
                    card_x = cards_left + idx * (card_w + gap)
                    intensity = str(month_info.get("hiring_intensity") or "moderate")
                    bar_color = intensity_colors.get(intensity, TEAL)

                    # Card background
                    _add_filled_rect(
                        slide, card_x, cards_top, card_w, card_h, WARM_WHITE
                    )

                    # Intensity color bar at top of card
                    _add_filled_rect(
                        slide, card_x, cards_top, card_w, Inches(0.08), bar_color
                    )

                    # Month name
                    _add_textbox(
                        slide,
                        card_x + Inches(0.08),
                        cards_top + Inches(0.15),
                        card_w - Inches(0.16),
                        Inches(0.25),
                        text=str(month_info.get("month_name") or ""),
                        font_size=11,
                        bold=True,
                        color=NAVY,
                    )

                    # Season + intensity
                    _add_textbox(
                        slide,
                        card_x + Inches(0.08),
                        cards_top + Inches(0.42),
                        card_w - Inches(0.16),
                        Inches(0.2),
                        text=f"{month_info.get('season', '')} | {intensity.replace('_', ' ').upper()}",
                        font_size=7,
                        bold=True,
                        color=bar_color if intensity != "low" else MUTED_TEXT,
                    )

                    # Budget weight
                    weight = month_info.get("budget_weight", 1.0)
                    _add_textbox(
                        slide,
                        card_x + Inches(0.08),
                        cards_top + Inches(0.65),
                        card_w - Inches(0.16),
                        Inches(0.2),
                        text=f"Budget weight: {weight:.1f}x",
                        font_size=7,
                        color=DARK_TEXT,
                    )

                    # Events list
                    events = month_info.get("key_events") or []
                    event_y = cards_top + Inches(0.9)
                    for ev in events[:3]:
                        _add_textbox(
                            slide,
                            card_x + Inches(0.08),
                            event_y,
                            card_w - Inches(0.16),
                            Inches(0.18),
                            text=f"- {ev}",
                            font_size=6,
                            color=MUTED_TEXT,
                        )
                        event_y += Inches(0.18)

                    # Recommendation
                    rec = str(month_info.get("recommendation") or "")
                    if rec:
                        _add_textbox(
                            slide,
                            card_x + Inches(0.08),
                            cards_top + Inches(1.7),
                            card_w - Inches(0.16),
                            Inches(0.6),
                            text=rec,
                            font_size=6,
                            italic=True,
                            color=DARK_TEXT,
                        )

            # ---- Budget Tiers Summary (bottom section) ----
            if budget_tiers and "error" not in budget_tiers:
                tier_top = Inches(5.3)
                _add_textbox(
                    slide,
                    Inches(0.55),
                    tier_top,
                    Inches(5),
                    Inches(0.25),
                    text="BUDGET TIER ALLOCATION",
                    font_size=10,
                    bold=True,
                    color=NAVY,
                )
                _add_filled_rect(
                    slide,
                    Inches(0.55),
                    tier_top + Inches(0.25),
                    Inches(1.5),
                    Inches(0.025),
                    TEAL,
                )

                tier_breakdown: dict = budget_tiers.get("tier_breakdown") or {}
                tier_x = Inches(0.55)
                tier_card_w = Inches(4.0)
                tier_card_h = Inches(1.0)

                for idx, (tier_key, tier_info) in enumerate(tier_breakdown.items()):
                    tx = tier_x + idx * (tier_card_w + Inches(0.1))
                    _add_filled_rect(
                        slide,
                        tx,
                        tier_top + Inches(0.4),
                        tier_card_w,
                        tier_card_h,
                        NAVY,
                    )

                    tier_label = tier_key.replace("_", " ").title()
                    amount = tier_info.get("amount", 0)
                    pct = tier_info.get("pct", 0)

                    _add_textbox(
                        slide,
                        tx + Inches(0.1),
                        tier_top + Inches(0.45),
                        tier_card_w - Inches(0.2),
                        Inches(0.35),
                        # S3: the plan's OWN budget-tier allocation -- localize.
                        text=_fmt_currency(amount),
                        font_size=16,
                        bold=True,
                        color=TEAL,
                        anchor=MSO_ANCHOR.MIDDLE,
                    )
                    _add_textbox(
                        slide,
                        tx + Inches(0.1),
                        tier_top + Inches(0.82),
                        tier_card_w - Inches(0.2),
                        Inches(0.25),
                        text=f"{tier_label} ({pct:.0f}%)",
                        font_size=8,
                        color=LIGHT_MUTED,
                        alignment=PP_ALIGN.CENTER,
                    )

            _add_footer(slide, today)

        except Exception as exc:
            logger.error(
                "Gold Standard slide (Activation Timeline) failed: %s",
                exc,
                exc_info=True,
            )


# ===================================================================
# SLIDE: Risk Analysis
# ===================================================================


def _build_slide_risk_analysis(prs: Presentation, data: Dict) -> None:
    """Build a Risk Analysis slide with budget, timing, channel, and competitive risks.

    Presents 4 risk categories with impact assessment and mitigation strategies,
    styled as a professional risk matrix for C-suite audiences.
    """
    try:
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        today = datetime.date.today().strftime("%B %d, %Y")

        client = data.get("client_name", "Client")
        budget_alloc = data.get("_budget_allocation", {})
        if not isinstance(budget_alloc, dict):
            budget_alloc = {}
        total_proj = budget_alloc.get("total_projected", {})
        if not isinstance(total_proj, dict):
            total_proj = {}
        channel_allocs = budget_alloc.get("channel_allocations", {})
        if not isinstance(channel_allocs, dict):
            channel_allocs = {}
        gold = data.get("_gold_standard") or {}
        competitor_map = gold.get("competitor_mapping") or {}
        budget_meta = budget_alloc.get("metadata", {})
        if not isinstance(budget_meta, dict):
            budget_meta = {}
        total_budget = budget_meta.get("total_budget") or 0

        # Background
        _add_filled_rect(
            slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, OFF_WHITE
        )
        _add_top_band(slide, "RISK ANALYSIS", today)

        _add_textbox(
            slide,
            Inches(0.55),
            Inches(0.92),
            Inches(12.2),
            Inches(0.45),
            text=f"Strategic risk assessment for {client}'s recruitment campaign",
            font_size=15,
            bold=True,
            color=NAVY,
        )

        # Build risk items
        risks: list[tuple[str, str, str, str]] = (
            []
        )  # (category, risk, impact, mitigation)

        # S48 FIX: per-channel sum for consistency
        proj_hires = sum(
            int(ch.get("projected_hires") or 0) for ch in channel_allocs.values()
        )
        if proj_hires == 0:
            proj_hires = total_proj.get("hires") or 0
        cph = (
            round(total_budget / max(proj_hires, 1), 2)
            if proj_hires > 0
            else (total_proj.get("cost_per_hire") or 0)
        )

        # 1. Budget risk
        if proj_hires > 0 and cph > 0 and total_budget > 0:
            hires_if_cpa_up = int(total_budget / (cph * 1.2))
            risks.append(
                (
                    "BUDGET",
                    "CPA Inflation Risk",
                    f"If CPA rises 20%, hires drop from {proj_hires:,.0f} to {hires_if_cpa_up:,.0f}",
                    "Build 10-15% budget contingency; diversify channels",
                )
            )
        else:
            risks.append(
                (
                    "BUDGET",
                    "Budget Uncertainty",
                    "Insufficient data for precise projection — actuals may vary 20-30%",
                    "Start with 2-week pilot; adjust allocation based on early CPA data",
                )
            )

        # 2. Market timing
        campaign_start = data.get("campaign_start_month") or 0
        if isinstance(campaign_start, int) and campaign_start in (4, 5, 6):
            risks.append(
                (
                    "TIMING",
                    "Q2 Competitive Peak",
                    "Q2 hiring is 15-20% more competitive due to fiscal year budget cycles",
                    "Front-load spend in weeks 1-4; lock niche channel inventory early",
                )
            )
        elif isinstance(campaign_start, int) and campaign_start in (1, 2, 3):
            risks.append(
                (
                    "TIMING",
                    "Q1 New Year Surge",
                    "25% more job seekers but 20% more employer competition in Q1",
                    "Capitalize on high candidate supply with aggressive apply-rate optimization",
                )
            )
        else:
            # copy:both#5: cite this plan's own 2-3 highest-intensity
            # months + actual events from the activation calendar instead
            # of the generic sentence that rendered byte-identically across
            # every bundle regardless of client/industry.
            _activation_cal = (gold.get("activation_calendar") or {}).get(
                "timeline"
            ) or []
            _high_months = [
                m
                for m in _activation_cal
                if isinstance(m, dict)
                and str(m.get("hiring_intensity") or "").lower()
                in ("high", "very_high")
            ][:3]
            if _high_months:
                _month_names = [str(m.get("month_name") or "") for m in _high_months]
                _events: list = []
                for m in _high_months:
                    _events.extend(str(e) for e in (m.get("key_events") or [])[:1])
                _events_str = f" ({', '.join(_events[:3])})" if _events else ""
                risks.append(
                    (
                        "TIMING",
                        "Seasonal Variations",
                        f"{', '.join(_month_names)} run high hiring intensity "
                        f"for this plan{_events_str}",
                        "Front-load budget into these months; monitor weekly "
                        "CPA trends and shift spend to high-performing periods",
                    )
                )
            else:
                risks.append(
                    (
                        "TIMING",
                        "Seasonal Variations",
                        "Hiring demand and candidate supply fluctuate throughout the year",
                        "Monitor weekly CPA trends; shift budget to high-performing periods",
                    )
                )

        # 3. Channel dependency
        if channel_allocs:
            sorted_ch = sorted(
                channel_allocs.items(),
                key=lambda x: (
                    x[1].get("percentage", 0) if isinstance(x[1], dict) else 0
                ),
                reverse=True,
            )
            top_2_names = [ch[0] for ch in sorted_ch[:2]]
            top_2_pct = sum(
                ch[1].get("percentage", 0) if isinstance(ch[1], dict) else 0
                for ch in sorted_ch[:2]
            )
            if top_2_pct > 55:
                # "Hires at risk" must be the ACTUAL sum of these named
                # channels' own projected_hires from _budget_allocation, not
                # a percentage-of-total-hires estimate -- the two can
                # diverge sharply (a shipped bundle once said "29 hires" here
                # while the named channels actually summed to 19).
                top_2_hires = sum(
                    int((channel_allocs.get(k) or {}).get("projected_hires") or 0)
                    for k in top_2_names
                )
                top_2_labels = [_fmt.channel_label(k) for k in top_2_names]
                risks.append(
                    (
                        "CHANNELS",
                        "Channel Concentration",
                        f"{top_2_pct:.0f}% of budget on {', '.join(top_2_labels)} — "
                        f"disruption could impact {_fmt.fmt_count(top_2_hires, 'hire')}",
                        "Diversify to 4+ channels; maintain backup channels on standby",
                    )
                )
            else:
                risks.append(
                    (
                        "CHANNELS",
                        "Channel Fragmentation",
                        "Budget spread across many channels may dilute impact",
                        "Consolidate on top 3-4 performers after 2-week pilot period",
                    )
                )
        else:
            risks.append(
                (
                    "CHANNELS",
                    "Channel Selection",
                    "Channel performance varies by role and location",
                    "A/B test top 3 channels in first 2 weeks before full commitment",
                )
            )

        # 4. Competitive pressure
        n_high_comp = sum(
            1
            for k, v in competitor_map.items()
            if not str(k).startswith("_")
            and isinstance(v, dict)
            and str(v.get("hiring_intensity") or "").lower() in ("high", "very_high")
        )
        if n_high_comp > 0:
            _n_high_comp_txt = _fmt.fmt_count(n_high_comp, "market")
            risks.append(
                (
                    "COMPETITION",
                    f"High Competition ({_n_high_comp_txt})",
                    f"Fortune 500+ companies hiring same roles in {_n_high_comp_txt}",
                    "Differentiate with employer brand; emphasize career growth and culture",
                )
            )
        else:
            # copy:both#5: name the plan's actual top competitor instead of
            # a generic, byte-identical-across-every-bundle sentence.
            # Brief-supplied competitors take precedence (same source
            # priority as the Competitive Landscape slide), falling back to
            # the gold-standard competitor mapping's top employer.
            _top_competitor = ""
            _direct_comps_r = data.get("competitors") or []
            if isinstance(_direct_comps_r, list) and _direct_comps_r:
                _first_c = _direct_comps_r[0]
                _top_competitor = (
                    str(_first_c).strip()
                    if isinstance(_first_c, str)
                    else (
                        str(_first_c.get("name") or "").strip()
                        if isinstance(_first_c, dict)
                        else ""
                    )
                )
            if not _top_competitor and competitor_map:
                for _ck, _cv in competitor_map.items():
                    if isinstance(_ck, str) and _ck.startswith("_"):
                        continue
                    if isinstance(_cv, dict) and _cv.get("top_employers"):
                        _top_competitor = _strip_competitor_tag(
                            str(_cv["top_employers"][0])
                        )
                        break
            if _top_competitor:
                risks.append(
                    (
                        "COMPETITION",
                        "Competitive Landscape",
                        f"{_top_competitor} and similar employers may increase "
                        f"hiring activity during the campaign period",
                        f"Monitor {_top_competitor}'s job posting volumes weekly; "
                        f"adjust messaging",
                    )
                )
            else:
                risks.append(
                    (
                        "COMPETITION",
                        "Competitive Landscape",
                        "Competitors may increase hiring activity during campaign period",
                        "Monitor competitor job posting volumes weekly; adjust messaging",
                    )
                )

        # Render risk cards in a 2x2 grid, sized to content (measure-then-
        # shrink, matching the O1 layout engine used on slides 2/5) instead
        # of a fixed 2.2in height that left large empty gaps between the
        # Impact and Mitigation lines whenever either was short
        # (visual:manpower#2 / visual:atria#3).
        card_w = Inches(6.0)
        card_w_in = 6.0
        gap = Inches(0.25)
        grid_left = Inches(0.55)
        _impact_top_in = 0.65
        _body_w_in = card_w_in - 0.4

        def _risk_impact_h_in(impact_text: str) -> float:
            n = _estimate_lines(f"Impact: {impact_text}", _body_w_in, 9)
            return n * (9 * 1.35) / 72.0

        def _risk_mitigation_h_in(mitigation_text: str) -> float:
            n = _estimate_lines(f"Mitigation: {mitigation_text}", _body_w_in, 9)
            return n * (9 * 1.35) / 72.0

        def _risk_card_h_in(impact_text: str, mitigation_text: str) -> float:
            mitigation_top_in = (
                _impact_top_in + _risk_impact_h_in(impact_text) + 0.18
            )
            return mitigation_top_in + _risk_mitigation_h_in(mitigation_text) + 0.22

        _risk_list = risks[:4]
        _row_h_in: list = []
        for row_start in range(0, len(_risk_list), 2):
            row_items = _risk_list[row_start : row_start + 2]
            row_h = max(
                _risk_card_h_in(impact, mitigation)
                for (_c, _t, impact, mitigation) in row_items
            )
            _row_h_in.append(max(1.5, min(row_h, 2.4)))

        # visual:manpower#4: vertically center the risk-card rows in the
        # space between the title and the footer rule instead of always
        # starting at a fixed y=1.5in -- short impact/mitigation text
        # previously left a large flat gap below the cards before the
        # footer. Never starts ABOVE the old y=1.5 anchor.
        _risk_avail_top_in = 1.5
        _risk_avail_bottom_in = 6.95  # footer rule sits at 7.12in
        _risk_content_h_in = sum(_row_h_in) + 0.25 * max(0, len(_row_h_in) - 1)
        _risk_start_y_in = _risk_avail_top_in + max(
            0.0,
            (_risk_avail_bottom_in - _risk_avail_top_in - _risk_content_h_in) / 2.0,
        )

        _row_y_in: list = []
        _cur_y_in = _risk_start_y_in
        for h_in in _row_h_in:
            _row_y_in.append(_cur_y_in)
            _cur_y_in += h_in + 0.25

        risk_colors = {
            "BUDGET": RED_ACCENT,
            "TIMING": AMBER,
            "CHANNELS": BLUE,
            "COMPETITION": TEAL,
        }

        for idx, (category, risk_title, impact, mitigation) in enumerate(_risk_list):
            col = idx % 2
            row_idx = idx // 2
            card_h_in = _row_h_in[row_idx]
            card_h = Inches(card_h_in)
            x = grid_left + col * (card_w + gap)
            y = Inches(_row_y_in[row_idx])

            # Card background
            _add_rounded_rect(slide, x, y, card_w, card_h, WHITE)
            accent_color = risk_colors.get(category, BLUE)
            _add_filled_rect(slide, x, y, Inches(0.06), card_h, accent_color)

            # Category badge
            _add_textbox(
                slide,
                x + Inches(0.2),
                y + Inches(0.08),
                Inches(1.5),
                Inches(0.25),
                text=category,
                font_size=8,
                bold=True,
                color=accent_color,
            )

            # Risk title
            _add_textbox(
                slide,
                x + Inches(0.2),
                y + Inches(0.32),
                card_w - Inches(0.4),
                Inches(0.28),
                text=risk_title,
                font_size=11,
                bold=True,
                color=NAVY,
            )

            # Impact
            impact_h_in = _risk_impact_h_in(impact)
            _add_textbox(
                slide,
                x + Inches(0.2),
                y + Inches(_impact_top_in),
                card_w - Inches(0.4),
                Inches(max(0.2, impact_h_in)),
                text=f"Impact: {impact}",
                font_size=9,
                color=DARK_TEXT,
            )

            # Mitigation
            mitigation_top_in = _impact_top_in + impact_h_in + 0.18
            _add_textbox(
                slide,
                x + Inches(0.2),
                y + Inches(mitigation_top_in),
                card_w - Inches(0.4),
                Inches(max(0.2, card_h_in - mitigation_top_in - 0.1)),
                text=f"Mitigation: {mitigation}",
                font_size=9,
                bold=False,
                color=MUTED_TEXT,
            )

        _add_footer(slide, today)

    except Exception as exc:
        import logging as _logging

        _logging.getLogger(__name__).error(
            "Risk Analysis slide failed: %s", exc, exc_info=True
        )


# ===================================================================
# DECK NARRATIVE SLIDES -- Joveo client deck 2026 parity
# (content from data/joveo_media_plan_deck_2026.json via _load_deck_kb)
# ===================================================================


def _build_slide_methodology(prs: Presentation, data: Dict, deck: Dict) -> None:
    """Build the 'Our Methodology' slide — Joveo's 6-step campaign management.

    Renders 6 numbered step cards (2 rows x 3 cols) with purple number chips,
    step name, AI-tool tag, and detail text. No-ops if the deck KB section
    is missing so generation never depends on the narrative content.
    """
    steps = ((deck.get("campaign_methodology") or {}).get("steps")) or []
    if not steps:
        return

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    today = datetime.date.today().strftime("%B %d, %Y")

    # Background + standard chrome
    _add_filled_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, OFF_WHITE)
    _add_top_band(slide, "OUR METHODOLOGY", today)

    _add_textbox(
        slide,
        Inches(0.55),
        Inches(0.92),
        Inches(12.2),
        Inches(0.45),
        text="Our Methodology — Joveo's 6-step campaign management on the MOJO platform",
        font_size=15,
        bold=True,
        color=NAVY,
    )

    # 2x3 grid of numbered step cards
    card_w = Inches(3.95)
    card_h = Inches(2.45)
    gap = Inches(0.2)
    grid_top = Inches(1.55)
    grid_left = Inches(0.55)

    for idx, step in enumerate(steps[:6]):
        if not isinstance(step, dict):
            continue
        col = idx % 3
        row_idx = idx // 3
        x = grid_left + col * (card_w + gap)
        y = grid_top + row_idx * (card_h + gap)

        accent = INV_ACCENTS[idx % len(INV_ACCENTS)]
        card = _add_rounded_rect(slide, x, y, card_w, card_h, WHITE)
        try:
            card.line.color.rgb = WARM_GRAY
            card.line.width = Pt(0.75)
        except Exception:
            pass
        # Colored top accent bar (Invisible card style)
        _add_filled_rect(slide, x, y, card_w, Inches(0.08), accent)

        # Number chip cycling the dataviz palette
        chip = Inches(0.42)
        _add_oval(slide, x + Inches(0.2), y + Inches(0.26), chip, chip, accent)
        _add_textbox(
            slide,
            x + Inches(0.2),
            y + Inches(0.26),
            chip,
            chip,
            text=str(step.get("n") or idx + 1),
            font_size=14,
            bold=True,
            color=WHITE,
            alignment=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )

        # Step name (Poppins bold, navy)
        _add_textbox(
            slide,
            x + Inches(0.74),
            y + Inches(0.24),
            card_w - Inches(0.92),
            Inches(0.66),
            text=str(step.get("name") or ""),
            font_size=11,
            bold=True,
            color=NAVY,
        )

        # AI-tool tag (small purple caps) -- step 1 has no AI tool
        ai_tag = str(step.get("ai") or "")
        if ai_tag:
            _add_textbox(
                slide,
                x + Inches(0.18),
                y + Inches(0.92),
                card_w - Inches(0.36),
                Inches(0.24),
                text=ai_tag.upper(),
                font_size=8,
                bold=True,
                color=BLUE,
            )

        # Detail text (Inter, muted)
        _box, _tf = _add_textbox(
            slide,
            x + Inches(0.18),
            y + Inches(1.2),
            card_w - Inches(0.36),
            Inches(1.12),
            text=str(step.get("detail") or ""),
            font_size=9,
            color=MUTED_TEXT,
        )
        _set_body_font(_tf)

    _add_footer(slide, today)


_PUSH_CHANNEL_KEYS = frozenset(
    {
        "programmatic_dsp",
        "programmatic",
        "niche_boards",
        "global_boards",
        "regional_boards",
        "job_board",
        "social_media",
        "social",
        "search_engine",
        "search",
        "display_retargeting",
        "display",
        "apac_regional",
        "emea_regional",
        "staffing_agency",
        "staffing",
        "events_jobfairs",
        "events",
        "lead_generation",
        "direct_hiring",
    }
)
_PULL_CHANNEL_KEYS = frozenset(
    {
        "employer_branding",
        "career_site",
        "referrals",
        "referral",
        "internal_mobility",
        "university",
        "email",
    }
)


def _push_pull_channel_split(data: Dict) -> "tuple[list, list]":
    """Split THIS plan's channel allocations (``_budget_allocation``) into
    Push (active/paid outreach: programmatic, boards, social, search) vs.
    Pull (brand magnetism: employer branding, career site, referrals,
    internal mobility) buckets, each a list of ``(label, dollar)`` sorted
    by spend descending.

    visual:manpower#2 / visual:atria#3: gives the Push Meets Pull slide
    real, plan-specific numbers instead of only generic prose.
    """
    budget_alloc = data.get("_budget_allocation", {})
    channel_allocs = (
        budget_alloc.get("channel_allocations", {})
        if isinstance(budget_alloc, dict)
        else {}
    )
    if not isinstance(channel_allocs, dict):
        return [], []
    push_items: list = []
    pull_items: list = []
    for ch_key, ch_data in channel_allocs.items():
        if not isinstance(ch_data, dict):
            continue
        dollar = ch_data.get("dollar_amount") or 0
        try:
            dollar = float(dollar)
        except (TypeError, ValueError):
            dollar = 0.0
        if dollar <= 0:
            continue
        label = _fmt.channel_label(str(ch_key))
        key_l = str(ch_key).strip().lower()
        if key_l in _PULL_CHANNEL_KEYS:
            pull_items.append((label, dollar))
        elif key_l in _PUSH_CHANNEL_KEYS:
            push_items.append((label, dollar))
        # Unclassified channel keys are skipped rather than guessed into
        # either bucket -- a wrong classification is worse than an omission.
    push_items.sort(key=lambda x: x[1], reverse=True)
    pull_items.sort(key=lambda x: x[1], reverse=True)
    return push_items, pull_items


def _push_pull_split_line(items: "list", total_budget: float) -> str:
    """One compact "This plan:" line of channel/$ for a Push or Pull bucket.

    visual:manpower#2 / consistency:manpower#1 / consistency:atria#1: every
    channel in ``items`` MUST be itemized -- a prior version capped the
    listed line items at ``items[:4]`` while ``total`` (and the printed
    "$X total") summed ALL of them, so a 5th/6th channel's dollars (e.g.
    Social Media) were silently folded into the printed total without ever
    being named. The textbox this renders into word-wraps, so a longer
    itemized list simply flows onto a second line instead of overflowing.
    """
    if not items:
        return ""
    total = sum(d for _, d in items)
    parts = [f"{label} {_fmt.fmt_money(d, compact=True)}" for label, d in items]
    pct = f" ({round(total / total_budget * 100)}% of budget)" if total_budget > 0 else ""
    return f"This plan: {' · '.join(parts)} — {_fmt.fmt_money(total, compact=True)} total{pct}"


def _build_slide_push_meets_pull(prs: Presentation, data: Dict, deck: Dict) -> None:
    """Build the 'Push Meets Pull' slide — active outreach vs. brand magnetism.

    Two rounded cards side-by-side: Push (lavender surface, purple pill)
    and Pull (blue-50 surface, teal-deep pill), with the deck summary line
    under the title band. Each card is sized to its actual content (measure-
    then-shrink, matching the O1 layout engine used on slides 2/5) instead
    of a fixed height, and carries a real "This plan:" channel/$ breakdown
    derived from ``_budget_allocation`` (visual:manpower#2 / visual:atria#3).
    No-ops if the deck KB section is missing.
    """
    pmp = deck.get("push_meets_pull") or {}
    push = pmp.get("push") or {}
    pull = pmp.get("pull") or {}
    if not push and not pull:
        return

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    today = datetime.date.today().strftime("%B %d, %Y")

    # Background + standard chrome
    _add_filled_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, OFF_WHITE)
    _add_top_band(slide, "PUSH MEETS PULL", today)

    _add_textbox(
        slide,
        Inches(0.55),
        Inches(0.92),
        Inches(12.2),
        Inches(0.45),
        text="Push Meets Pull — a dual-engine sourcing strategy",
        font_size=15,
        bold=True,
        color=NAVY,
    )

    # Summary line under the title band (Inter, muted)
    summary = str(pmp.get("summary") or "")
    if summary:
        _box, _tf = _add_textbox(
            slide,
            Inches(0.55),
            Inches(1.4),
            Inches(12.2),
            Inches(0.35),
            text=summary,
            font_size=11,
            color=MUTED_TEXT,
        )
        _set_body_font(_tf)

    push_split, pull_split = _push_pull_channel_split(data)
    _total_budget = sum(d for _, d in push_split) + sum(d for _, d in pull_split)

    # Two rounded cards side-by-side, measure-then-shrink to content.
    card_w_in = 6.0
    _detail_w_in = card_w_in - 0.8
    _pill_block_in = 0.4 + 0.55 + 0.25  # pill offset + pill height + gap to body
    _bottom_pad_in = 0.3

    cards = [
        (Inches(0.55), push, LAVENDER_100, BLUE, push_split),
        (Inches(6.8), pull, BLUE_50, AMBER, pull_split),
    ]

    # Measure BOTH cards first so they can share one common height (keeps the
    # side-by-side layout visually balanced instead of two different heights).
    _measured_h: list = []
    for _card_x, section, _surface, _pill_color, _split in cards:
        if not isinstance(section, dict) or not section:
            _measured_h.append(0.0)
            continue
        detail_text = str(section.get("detail") or "")
        split_line = _push_pull_split_line(_split, _total_budget)
        n_detail_lines = _estimate_lines(detail_text, _detail_w_in, 11) if detail_text else 0
        n_split_lines = (
            _estimate_lines(split_line, _detail_w_in, 9, char_em=0.5) if split_line else 0
        )
        detail_h = n_detail_lines * (11 * 1.35) / 72.0
        split_h = (n_split_lines * (9 * 1.35) / 72.0 + 0.15) if split_line else 0.0
        _measured_h.append(_pill_block_in + detail_h + split_h + _bottom_pad_in)

    card_h_in = max([2.2] + _measured_h)
    card_h_in = min(card_h_in, 4.6)  # never exceed the old fixed max
    card_h = Inches(card_h_in)

    # visual:manpower#4: vertically center the card row in the space between
    # the summary line and the footer rule, instead of always starting at a
    # fixed y=2.0in -- short content (few channels, terse detail copy)
    # previously left 35-45% flat dead space below the cards before the
    # footer. Centering distributes the slack above/below instead of only
    # below, and never moves the cards ABOVE the old y=2.0 starting point.
    _avail_top_in = 2.0
    _avail_bottom_in = 6.95  # footer rule sits at 7.12in
    card_top_in = _avail_top_in + max(
        0.0, (_avail_bottom_in - _avail_top_in - card_h_in) / 2.0
    )
    card_top = Inches(card_top_in)

    for card_x, section, surface, pill_color, split in cards:
        if not isinstance(section, dict) or not section:
            continue
        _add_rounded_rect(slide, card_x, card_top, Inches(card_w_in), card_h, surface)

        # Pill title
        _add_rounded_rect(
            slide,
            card_x + Inches(0.4),
            card_top + Inches(0.4),
            Inches(4.8),
            Inches(0.55),
            pill_color,
        )
        _add_textbox(
            slide,
            card_x + Inches(0.4),
            card_top + Inches(0.4),
            Inches(4.8),
            Inches(0.55),
            text=str(section.get("name") or ""),
            font_size=12,
            bold=True,
            color=WHITE,
            alignment=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )

        # Detail text (Inter)
        detail_text = str(section.get("detail") or "")
        _detail_top = card_top + Inches(_pill_block_in)
        _box, _tf = _add_textbox(
            slide,
            card_x + Inches(0.4),
            _detail_top,
            Inches(_detail_w_in),
            card_h - Inches(_pill_block_in) - Inches(_bottom_pad_in),
            text=detail_text,
            font_size=11,
            color=DARK_TEXT,
        )
        _set_body_font(_tf)

        # This plan's real $ split for this bucket (visual:manpower#2 /
        # visual:atria#3): plan-specific numbers, not more prose.
        split_line = _push_pull_split_line(split, _total_budget)
        if split_line:
            n_detail_lines = (
                _estimate_lines(detail_text, _detail_w_in, 11) if detail_text else 0
            )
            _split_top = _detail_top + Inches(
                max(0.3, n_detail_lines * (11 * 1.35) / 72.0) + 0.1
            )
            _sbox, _stf = _add_textbox(
                slide,
                card_x + Inches(0.4),
                _split_top,
                Inches(_detail_w_in),
                card_h - (_split_top - card_top) - Inches(0.1),
                text=split_line,
                font_size=9,
                bold=True,
                color=pill_color,
            )
            _set_body_font(_stf)

    _add_footer(slide, today)


def _role_breakdown_median_salary(
    gold: Dict[str, Any], title: str
) -> "tuple[Optional[float], bool]":
    """Look up ``title``'s median salary from the SAME
    ``_gold_standard.city_level_data[*].per_role_salary`` structure the
    workbook's Quality Intelligence "Salary Intelligence" table renders
    (excel_v2._build_sheet_quality_intelligence), instead of
    ``_enriched.salary_data`` -- which is fully absent whenever live
    enrichment is offline, even though gold_standard's per-role figures are
    already computed and sitting in the same bundle (consistency:atria#2).

    Multi-location plans get one ``per_role_salary`` entry per city; this
    averages the role's median across every city that has a matching entry
    (equal weight per market) rather than arbitrarily picking one city.
    Returns ``(median, is_estimated)`` -- ``is_estimated`` is True if ANY
    contributing city's entry has ``confidence == "estimated"`` (a
    tier-scaled fallback rather than a matched industry benchmark), so the
    caller can append the same "(est.)" honesty marker the workbook uses.
    Returns ``(None, False)`` when no city has data for this role.
    """
    city_level = gold.get("city_level_data") if isinstance(gold, dict) else None
    if not isinstance(city_level, dict) or not city_level:
        return None, False
    medians: list = []
    any_estimated = False
    for city_info in city_level.values():
        if not isinstance(city_info, dict):
            continue
        role_salary = city_info.get("per_role_salary") or {}
        if not isinstance(role_salary, dict):
            continue
        sal = role_salary.get(title)
        if not isinstance(sal, dict):
            continue
        median = sal.get("median")
        try:
            median = float(median)
        except (TypeError, ValueError):
            continue
        medians.append(median)
        if sal.get("confidence") == "estimated":
            any_estimated = True
    if not medians:
        return None, False
    return sum(medians) / len(medians), any_estimated


def _build_slide_role_breakdown(prs: Presentation, data: Dict) -> None:
    """Compact role-by-role table: tier, est. median salary, and channel
    emphasis per role -- sourced from ``_gold_standard.difficulty_framework``
    and ``_gold_standard.city_level_data[*].per_role_salary`` (the SAME
    per-role salary data the workbook's Quality Intelligence sheet renders;
    see :func:`_role_breakdown_median_salary` -- consistency:atria#2).

    Shown only when the plan names >=4 roles AND gold-standard role data
    exists for at least 4 of them; no-ops cleanly otherwise. Fills the slot
    left empty when ``_build_slide_cpa_reference`` drops (this plan's
    industry has no KB cpa_reference content) -- visual:atria#4: a 10-role
    plan previously had NO role-level view anywhere in the deck.
    """
    roles = data.get("target_roles") or data.get("roles") or []
    if isinstance(roles, str):
        roles = [roles]
    # Roles may be plain strings or {"title": ..., "count": ..., "tier": ...}
    # dicts (the shape app.py's own role normalization produces) -- str(r)
    # on a dict would render its Python repr as the "role title".
    role_titles = []
    for r in roles or []:
        if isinstance(r, dict):
            _rt = str(r.get("title") or r.get("role") or "").strip()
        else:
            _rt = str(r).strip()
        if _rt:
            role_titles.append(_rt)
    if len(role_titles) < 4:
        return

    gold = data.get("_gold_standard") or {}
    difficulty = gold.get("difficulty_framework") or []
    if not isinstance(difficulty, list) or not difficulty:
        return
    by_title = {
        str(d.get("role_title") or "").strip().lower(): d
        for d in difficulty
        if isinstance(d, dict)
    }
    if not by_title:
        return

    rows: list = []
    for title in role_titles:
        d = by_title.get(title.lower())
        if not d:
            continue
        tier = _fmt.smart_title(str(d.get("seniority_level") or ""))
        emphasis = _fmt.channel_label(
            str(d.get("channel_emphasis") or "")
        ) or _fmt.smart_title(str(d.get("channel_emphasis") or ""))
        # consistency:atria#2: read the SAME per_role_salary data the
        # workbook renders, not the (often-empty) _enriched.salary_data.
        median, is_estimated = _role_breakdown_median_salary(gold, title)
        # strategy:atria#5: surface the same per-role Difficulty
        # (complexity_score) and Budget Weight the workbook's Role
        # Difficulty Classification table carries, so a 10-role plan with a
        # 20x salary/difficulty spread has real per-role substance here --
        # not just tier/salary/channel with zero hire/budget signal.
        rows.append(
            {
                "title": title,
                "tier": tier,
                "median": median,
                "is_estimated": is_estimated,
                "difficulty": d.get("complexity_score"),
                "budget_weight": d.get("budget_weight"),
                "emphasis": emphasis,
            }
        )
    if len(rows) < 4:
        return

    # visual:atria#2: when an ESTIMATED (fallback) salary happens to be
    # byte-identical to another role's displayed salary -- e.g. an
    # entry-level "Sales" role's fallback landing on the exact same $104K as
    # a mid-tier Nurse -- that coincidence is itself the tell that the
    # fallback defaulted to the wrong occupation code. Flag it explicitly
    # rather than silently rendering the same implausible number twice.
    _median_counts: Dict[float, int] = {}
    for r in rows:
        if r["median"]:
            _median_counts[r["median"]] = _median_counts.get(r["median"], 0) + 1
    for r in rows:
        if r["median"]:
            salary_str = _format_salary(r["median"])
            if r["is_estimated"] and salary_str:
                if _median_counts.get(r["median"], 0) > 1:
                    salary_str += " (est., shared band)"
                else:
                    salary_str += " (est.)"
        else:
            salary_str = "--"
        r["salary_str"] = salary_str
        r["difficulty_str"] = (
            f"{float(r['difficulty']):.1f}/10"
            if isinstance(r["difficulty"], (int, float))
            else "--"
        )
        r["budget_weight_str"] = (
            f"{float(r['budget_weight']):.1f}x"
            if isinstance(r["budget_weight"], (int, float))
            else "--"
        )

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    today = datetime.date.today().strftime("%B %d, %Y")
    client = data.get("client_name", "Client")

    _add_filled_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, OFF_WHITE)
    _add_top_band(slide, "ROLE BREAKDOWN", today)
    _add_textbox(
        slide,
        Inches(0.55),
        Inches(0.92),
        Inches(12.2),
        Inches(0.45),
        text=(
            f"Tier, salary band, and channel emphasis across {client}'s "
            f"{len(rows)} target roles"
        ),
        font_size=15,
        bold=True,
        color=NAVY,
    )

    table_left = Inches(0.55)
    # strategy:atria#5: widened from 4 to 6 columns (added Difficulty and
    # Budget Weight, sourced from the same _gold_standard difficulty_
    # framework rows the workbook's Role Difficulty Classification table
    # renders) so a 10-role plan carries real per-role substance here.
    col_widths = [
        Inches(3.1),
        Inches(1.2),
        Inches(2.15),
        Inches(1.25),
        Inches(1.5),
        Inches(2.8),
    ]
    headers = [
        "Role",
        "Tier",
        "Est. Median Salary",
        "Difficulty",
        "Budget Weight",
        "Channel Emphasis",
    ]
    header_top = Inches(1.7)
    header_h = Inches(0.42)
    row_h = Inches(0.42)

    _add_filled_rect(
        slide, table_left, header_top, sum(col_widths, Inches(0)), header_h, NAVY
    )
    cx = table_left
    for header, cw in zip(headers, col_widths):
        _add_textbox(
            slide,
            cx + Inches(0.15),
            header_top,
            cw - Inches(0.3),
            header_h,
            text=header,
            font_size=10,
            bold=True,
            color=WHITE,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        cx += cw

    _max_rows = min(len(rows), 12)
    for idx, r in enumerate(rows[:_max_rows]):
        y = header_top + header_h + idx * row_h
        _add_filled_rect(
            slide,
            table_left,
            y,
            sum(col_widths, Inches(0)),
            row_h,
            LAVENDER_50 if idx % 2 == 0 else WHITE,
        )
        cells = [
            (_trunc_word(r["title"], 40), 9, True, NAVY),
            (r["tier"], 9, False, DARK_TEXT),
            (r["salary_str"], 9, False, DARK_TEXT),
            (r["difficulty_str"], 9, False, DARK_TEXT),
            (r["budget_weight_str"], 9, False, DARK_TEXT),
            (r["emphasis"], 9, False, MUTED_TEXT),
        ]
        cx = table_left
        for (cell_text, fsize, fbold, fcolor), cw in zip(cells, col_widths):
            _add_textbox(
                slide,
                cx + Inches(0.15),
                y,
                cw - Inches(0.3),
                row_h,
                text=cell_text,
                font_size=fsize,
                bold=fbold,
                color=fcolor,
                anchor=MSO_ANCHOR.MIDDLE,
            )
            cx += cw

    _add_footer(slide, today)


def _build_slide_cpa_reference(prs: Presentation, data: Dict, deck: Dict) -> None:
    """Build the 'CPA Reference' slide — baseline CPA ranges by role category.

    Shape-built table: indigo header row with white Poppins text, columns
    Role Category | Est. CPA Range | Benchmark Basis, zebra-striped rows.
    KB v2.0 schema: ``cpa_reference[<industry_key>]`` (industry-keyed; only
    ``'ai_training'`` is populated). No-ops (drops the slide entirely, no
    empty placeholder) when this plan's key has no KB content.
    """
    cpa_by_key = deck.get("cpa_reference") or {}
    key = _kb_content_key(data)
    cpa = cpa_by_key.get(key) if isinstance(cpa_by_key, dict) else None
    if not isinstance(cpa, dict):
        return
    roles = cpa.get("roles") or []
    if not roles:
        return

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    today = datetime.date.today().strftime("%B %d, %Y")

    # Background + standard chrome
    _add_filled_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, OFF_WHITE)
    _add_top_band(slide, "CPA REFERENCE", today)

    _add_textbox(
        slide,
        Inches(0.55),
        Inches(0.92),
        Inches(12.2),
        Inches(0.45),
        text="CPA Reference — baseline cost-per-application by role category",
        font_size=15,
        bold=True,
        color=NAVY,
    )

    # KB calibration note (Inter, muted)
    desc = str(cpa.get("description") or "")
    if desc:
        _box, _tf = _add_textbox(
            slide,
            Inches(0.55),
            Inches(1.38),
            Inches(12.2),
            Inches(0.3),
            text=desc,
            font_size=9,
            color=MUTED_TEXT,
        )
        _set_body_font(_tf)

    # Shape-built table: header + zebra-striped rows
    # S3: this KB table (joveo_media_plan_deck_2026.json cpa_reference) is
    # explicitly "currency": "USD" -- a fixed US/global-remote benchmark
    # constant, not the current plan's own figure. Mark it directly on the
    # column header (inline on the figure), not only in the description
    # caption above, per the no-bare-$ rule.
    table_left = Inches(0.55)
    col_widths = [Inches(4.6), Inches(2.4), Inches(5.2)]
    headers = ["Role Category", "Est. CPA Range (USD)", "Benchmark Basis"]
    header_top = Inches(1.8)
    header_h = Inches(0.42)
    row_h = Inches(0.58)

    # Header row -- indigo with white Poppins text
    _add_filled_rect(slide, table_left, header_top, sum(col_widths, Inches(0)), header_h, NAVY)
    cx = table_left
    for header, cw in zip(headers, col_widths):
        _add_textbox(
            slide,
            cx + Inches(0.15),
            header_top,
            cw - Inches(0.3),
            header_h,
            text=header,
            font_size=10,
            bold=True,
            color=WHITE,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        cx += cw

    for idx, role in enumerate(roles[:7]):
        if not isinstance(role, dict):
            continue
        y = header_top + header_h + idx * row_h

        # Zebra striping on alternating rows
        _add_filled_rect(
            slide,
            table_left,
            y,
            sum(col_widths, Inches(0)),
            row_h,
            LAVENDER_50 if idx % 2 == 0 else WHITE,
        )

        low = role.get("cpa_low")
        high = role.get("cpa_high")
        if low is not None and high is not None:
            cpa_range = f"${low} – ${high}"
            # S3: this KB table is explicitly "currency": "USD" -- mark the
            # figure itself (not just the column header) on a non-USD plan.
            if _get_active_currency() != "USD":
                cpa_range = _mark_usd(cpa_range)
        else:
            cpa_range = "--"
        cells = [
            (str(role.get("category") or ""), 10, True, NAVY),
            (cpa_range, 10, True, BLUE),
            (str(role.get("basis") or ""), 9, False, MUTED_TEXT),
        ]
        cx = table_left
        for (cell_text, fsize, fbold, fcolor), cw in zip(cells, col_widths):
            _box, _tf = _add_textbox(
                slide,
                cx + Inches(0.15),
                y,
                cw - Inches(0.3),
                row_h,
                text=cell_text,
                font_size=fsize,
                bold=fbold,
                color=fcolor,
                anchor=MSO_ANCHOR.MIDDLE,
            )
            if not fbold:  # basis column reads as body copy
                _set_body_font(_tf)
            cx += cw

    _add_footer(slide, today)


def _build_slide_why_joveo(prs: Presentation, data: Dict, deck: Dict) -> None:
    """Build 'The Joveo Advantage' -- four differentiator cards (Invisible style).

    Uses client-appropriate, universal Joveo value props (not the Invisible-
    specific KB narrative) so the slide reads correctly for any client. Pure
    Invisible accent-bar card layout + dark takeaway callout.
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    today = datetime.date.today().strftime("%B %d, %Y")
    client = data.get("client_name", "Client")
    industry_label = data.get("industry_label") or data.get("industry") or "your"

    _add_top_band(slide, "The Joveo Advantage", today)
    _add_textbox(
        slide, Inches(0.63), Inches(0.92), Inches(12.2), Inches(0.45),
        text=f"Why Joveo is the right partner for {client}'s {industry_label} hiring at scale",
        font_size=13, italic=True, color=BLUE,
    )

    cards = [
        (
            "AI-Powered Optimization",
            "Real-time bid & budget optimization at the individual-job level — spend "
            "shifts autonomously from easy-to-fill to hard-to-fill roles.",
        ),
        (
            "Global, Multi-Language Reach",
            "Proven playbooks across 50+ countries — databases, DSPs, social, freelance "
            "platforms and local job boards, in local languages.",
        ),
        (
            "Job Content Optimization",
            "A JD scoring engine plus automated job creation produce high-converting, "
            "SEO-optimized listings that lift apply rates.",
        ),
        (
            "Managed, Transparent Service",
            "Joveo runs every channel and database for you — transparent pay-per-click / "
            "apply pricing, scaled or shifted weekly.",
        ),
    ]
    cw = Inches(6.0)
    # visual:manpower#2 / visual:atria#3: size the card to content (measure-
    # then-shrink, matching the O1 layout engine on slides 2/5) instead of a
    # fixed 2.05in height that left large dead space below two sentences of
    # copy. ``_inv_card`` places its body text at y+0.64in with a
    # 0.24in horizontal pad on each side.
    _joveo_body_w_in = 6.0 - 0.24 * 2

    def _joveo_card_h_in(body_text: str) -> float:
        n = _estimate_lines(body_text, _joveo_body_w_in, 11)
        body_h_in = n * (11 * 1.35) / 72.0
        return 0.64 + body_h_in + 0.16

    ch_in = max(1.35, min(2.05, max(_joveo_card_h_in(b) for _, b in cards)))
    ch = Inches(ch_in)
    gap_x = Inches(0.3)
    gap_y = Inches(0.25)
    left = Inches(0.65)
    top = Inches(1.62)
    for i, (header, body) in enumerate(cards):
        col = i % 2
        row = i // 2
        x = left + col * (cw + gap_x)
        y = top + Emu(int(row * (int(ch) + int(gap_y))))
        _inv_card(
            slide, x, y, cw, ch, header, body,
            accent=INV_ACCENTS[i % len(INV_ACCENTS)], body_size=11,
        )

    # visual:manpower#4: pull the closing callout band up to sit right below
    # the card grid instead of always at a fixed y=6.5in -- when the cards'
    # measured height is near its 1.35in floor (short body copy), leaving
    # the callout pinned at 6.5 created a wide flat gap between the cards
    # and the banner. Never moves the callout LOWER than the old default
    # (min(...) keeps it a comfortable distance above the footer rule at
    # 7.12in when the cards are tall).
    _grid_bottom_in = 1.62 + 2 * ch_in + 0.25
    _callout_top_in = min(6.5, _grid_bottom_in + 0.35)
    _inv_callout(
        slide,
        f"One managed partner, one tracked funnel — built to scale {client}'s hiring "
        f"with a predictable cost-per-hire.",
        top=_callout_top_in,
    )
    _add_footer(slide, today)


def _is_unbounded_duration(duration: Any) -> bool:
    """True when ``duration`` represents an unbounded/open-ended campaign --
    the wizard's "Ongoing" duration option, any case variant of it, or an
    empty/not-specified value -- rather than a fixed length like "6 months".

    Mirrors ``excel_v2._is_unbounded_duration`` exactly (kept as a local
    copy rather than a cross-module import -- this codebase's established
    pattern for small helpers shared between excel_v2.py and
    ppt_generator.py, e.g. ``_proper_client_name``/``_TITLE_ACRONYMS``) so
    the deck's Next Steps slide never renders the nonsensical "over
    Ongoing" phrasing the workbook's executive summary used to produce
    (prod defect: "Finalize weekly budget (25000 over Ongoing) and success
    metrics").
    """
    s = str(duration or "").strip().lower()
    return s in ("", "ongoing", "unbounded", "not specified", "tbd", "n/a")


def _interpolate_next_steps(steps: List[Any], data: Dict) -> List[str]:
    """Splice this plan's own facts into the KB's industry-agnostic Next
    Steps template (``data/joveo_media_plan_deck_2026.json``).

    copy:both#4: the 5 Next Steps bullets came verbatim from the KB, so a
    logistics client and a healthcare client got byte-identical text with
    no client name, role, location, or dollar figure anywhere on the slide.
    This matches each KB bullet by its stable keyword phrase and splices in
    the client display name, top roles, location count/names, budget, and
    duration -- never fabricating a number the plan doesn't have. Any
    bullet whose wording doesn't match one of the known KB phrases (e.g. a
    future KB edit) is passed through unchanged rather than guessed at.
    """
    if not steps:
        return []

    client = _proper_client_name(str(data.get("client_name") or "").strip()) or ""

    roles = data.get("target_roles") or data.get("roles") or []
    if isinstance(roles, str):
        roles = [roles]
    role_names: list = []
    for r in roles or []:
        if isinstance(r, dict):
            rt = str(r.get("title") or r.get("role") or "").strip()
        else:
            rt = str(r).strip()
        if rt:
            role_names.append(rt)
    roles_phrase = _cap_roles_for_headline(role_names, cap=2) if role_names else ""

    locations = data.get("locations") or []
    if isinstance(locations, str):
        locations = [locations]
    loc_names = [str(loc).strip() for loc in locations if str(loc).strip()]
    if len(loc_names) == 1:
        loc_phrase = f"{loc_names[0]} market"
    elif len(loc_names) == 2:
        loc_phrase = f"{' and '.join(loc_names)} markets"
    elif len(loc_names) > 2:
        loc_phrase = (
            f"{len(loc_names)} markets ({', '.join(loc_names[:2])} "
            f"and {len(loc_names) - 2} more)"
        )
    else:
        loc_phrase = ""

    # copy:both#4 follow-up / S95-style defect: ``data["budget"]`` can arrive
    # either as a pre-formatted display string ("$150,000", from the wizard)
    # or a raw number (25000, from a JSON brief) -- ``_parse_budget_number``
    # (shared_utils.parse_budget_display) normalizes either shape to a float
    # so it can be rendered through the SAME compact formatter used
    # everywhere else in this deck, instead of a raw "25000" leaking onto
    # the slide verbatim. An unbounded ("Ongoing") duration gets its own
    # natural phrasing (mirroring excel_v2._build_deterministic_executive_
    # summary / _is_unbounded_duration above) instead of the grammatically
    # broken "over Ongoing" -- worded slightly differently for the
    # parenthetical vs. em-dash bullet so each reads as a single clause.
    budget_num = _parse_budget_number(data.get("budget"))
    budget_fmt = _fmt.fmt_money(budget_num, compact=True) if budget_num else ""
    duration = str(data.get("campaign_duration") or "").strip()
    unbounded = bool(duration) and _is_unbounded_duration(duration)

    def _budget_duration_phrase(dash_style: bool) -> str:
        if budget_fmt and duration:
            if unbounded:
                return (
                    f"{budget_fmt} on an ongoing basis"
                    if dash_style
                    else f"{budget_fmt}, ongoing"
                )
            return f"{budget_fmt} over {duration}"
        if budget_fmt:
            return budget_fmt
        if duration:
            if unbounded:
                return "on an ongoing basis" if dash_style else "ongoing"
            return f"over {duration}"
        return ""

    budget_duration_paren = _budget_duration_phrase(dash_style=False)
    budget_duration_dash = _budget_duration_phrase(dash_style=True)

    out: List[str] = []
    for step in steps:
        text = str(step)
        low = text.lower()
        if "align on scope" in low and roles_phrase and loc_phrase:
            text = f"Align on scope: confirm the {roles_phrase} priority and the {loc_phrase}"
        elif "finalize weekly budget" in low and budget_duration_paren:
            text = f"Finalize weekly budget ({budget_duration_paren}) and success metrics"
        elif "integrate the client job feed" in low and client:
            text = f"Integrate {client}'s job feed with the Joveo platform"
        elif "launch campaign within" in low and budget_duration_dash:
            text = f"Launch within 2 business weeks of feed integration — {budget_duration_dash}"
        elif "30-day review" in low and client:
            text = f"30-day review: CPA actuals, pipeline quality, {client}'s scale-up recommendation"
        out.append(text)
    return out


def _interpolate_timeline_bullets(
    phases: List[Dict[str, Any]], data: Dict, channels: Dict
) -> None:
    """copy:both#2: splice this plan's own facts into the Implementation
    Timeline's phase action items, which used to be byte-for-byte identical
    boilerplate regardless of client/industry/duration -- an 18-month,
    10-role, $300K senior-living ramp read identically to a 6-month,
    1-role, $150K trucking launch. Mutates ``phases`` (each a dict with a
    "bullets" list) in place; matches each bullet by its stable keyword
    phrase, mirroring :func:`_interpolate_next_steps`. Week counts are
    already plan-specific (derived from ``campaign_weeks`` by the caller);
    this only touches the bullet TEXT.
    """
    client = _proper_client_name(str(data.get("client_name") or "").strip())
    sorted_ch = sorted(
        channels.values(), key=lambda c: c.get("pct", 0), reverse=True
    )
    top_channel_names = [ch["label"] for ch in sorted_ch[:2] if ch.get("label")]
    top_channels_phrase = (
        " and ".join(top_channel_names) if top_channel_names else ""
    )
    for ph in phases:
        bullets = ph.get("bullets") or []
        new_bullets = []
        for b in bullets:
            text = str(b)
            low = text.lower()
            if "campaign setup" in low and "publisher activation" in low and client:
                text = f"Campaign setup & publisher activation for {client}"
            elif low.strip() == "scale top performers" and top_channels_phrase:
                text = f"Scale top performers: {top_channels_phrase}"
            new_bullets.append(text)
        ph["bullets"] = new_bullets


def _build_slide_next_steps_only(prs: Presentation, next_steps: List[Any]) -> None:
    """Full-canvas 'Next Steps' slide, used when this plan's industry has no
    KB case-study content (see :func:`_kb_content_key`). Gives the numbered
    steps the whole canvas -- deliberately renders NO case-study header,
    stat tiles, or Challenges/Solution columns, since those would either be
    empty or would require content this plan's industry doesn't have.
    No-ops if there are no next steps either.
    """
    if not next_steps:
        return

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    today = datetime.date.today().strftime("%B %d, %Y")

    _add_filled_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, OFF_WHITE)
    _add_top_band(slide, "NEXT STEPS", today)
    _add_textbox(
        slide,
        Inches(0.55),
        Inches(0.92),
        Inches(12.2),
        Inches(0.45),
        text="Recommended path forward",
        font_size=15,
        bold=True,
        color=NAVY,
    )

    steps = list(next_steps)[:6]
    n = len(steps)
    top_in = 1.75
    bottom_in = 6.9
    row_h_in = (bottom_in - top_in) / n
    for idx, step in enumerate(steps):
        y = Inches(top_in + idx * row_h_in)
        row_h = Inches(row_h_in - 0.14)
        _add_rounded_rect(
            slide,
            Inches(0.55),
            y,
            Inches(12.25),
            row_h,
            LAVENDER_50 if idx % 2 == 0 else WHITE,
        )
        chip = Inches(0.44)
        chip_y = y + Emu(int((int(row_h) - int(chip)) / 2))
        _add_oval(slide, Inches(0.85), chip_y, chip, chip, BLUE)
        _add_textbox(
            slide,
            Inches(0.85),
            chip_y,
            chip,
            chip,
            text=str(idx + 1),
            font_size=13,
            bold=True,
            color=WHITE,
            alignment=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        _box, _tf = _add_textbox(
            slide,
            Inches(1.55),
            y,
            Inches(10.9),
            row_h,
            text=str(step),
            font_size=13,
            color=DARK_TEXT,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        _set_body_font(_tf)

    _add_footer(slide, today)


def _build_slide_case_study_next_steps(
    prs: Presentation, data: Dict, deck: Dict
) -> None:
    """Build the 'Case Study & Next Steps' closing slide.

    Top: 3 KPI stat blocks from case_study.results. Middle: Challenges vs.
    The Joveo Solution bullet columns. Bottom: numbered Next Steps strip.

    KB v2.0 schema: ``case_study[<industry_key>]`` (industry-keyed; only
    ``'ai_training'`` is populated). When this plan's key has no KB case
    study, this NEVER fabricates one -- it falls through to
    :func:`_build_slide_next_steps_only`, a full-canvas Next Steps layout
    with no case-study header/tiles/stat row. No-ops entirely if there is
    also no ``next_steps`` content.
    """
    case_by_key = deck.get("case_study") or {}
    key = _kb_content_key(data)
    case = case_by_key.get(key) if isinstance(case_by_key, dict) else None
    # copy:both#4: splice this plan's own client name/roles/locations/
    # budget/duration into the KB's generic Next Steps template.
    next_steps = _interpolate_next_steps(deck.get("next_steps") or [], data)
    if not isinstance(case, dict) or not case:
        _build_slide_next_steps_only(prs, next_steps)
        return

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    today = datetime.date.today().strftime("%B %d, %Y")

    # Background + standard chrome
    _add_filled_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, OFF_WHITE)
    _add_top_band(slide, "CASE STUDY & NEXT STEPS", today)

    _add_textbox(
        slide,
        Inches(0.55),
        Inches(0.92),
        Inches(12.2),
        Inches(0.45),
        text=f"Case study: {case.get('title') or 'Proven results at scale'}",
        font_size=15,
        bold=True,
        color=NAVY,
    )

    # Top: 3 KPI stat blocks (value big purple, detail small muted)
    results = case.get("results") or []
    block_w = Inches(3.95)
    block_h = Inches(1.5)
    block_top = Inches(1.45)
    gap = Inches(0.2)
    for idx, res in enumerate(results[:3]):
        if not isinstance(res, dict):
            continue
        x = Inches(0.55) + idx * (block_w + gap)
        accent = INV_ACCENTS[idx % len(INV_ACCENTS)]
        block_card = _add_rounded_rect(slide, x, block_top, block_w, block_h, WHITE)
        try:
            block_card.line.color.rgb = WARM_GRAY
            block_card.line.width = Pt(0.75)
        except Exception:
            pass
        _add_filled_rect(slide, x, block_top, block_w, Inches(0.08), accent)
        _add_textbox(
            slide,
            x + Inches(0.2),
            block_top + Inches(0.18),
            block_w - Inches(0.4),
            Inches(0.5),
            text=str(res.get("value") or ""),
            font_size=28,
            bold=True,
            color=accent,
        )
        _add_textbox(
            slide,
            x + Inches(0.2),
            block_top + Inches(0.72),
            block_w - Inches(0.4),
            Inches(0.24),
            text=str(res.get("metric") or "").upper(),
            font_size=8,
            bold=True,
            color=NAVY,
        )
        _box, _tf = _add_textbox(
            slide,
            x + Inches(0.2),
            block_top + Inches(0.98),
            block_w - Inches(0.4),
            Inches(0.46),
            text=str(res.get("detail") or ""),
            font_size=8,
            color=MUTED_TEXT,
        )
        _set_body_font(_tf)

    # Middle: Challenges vs. The Joveo Solution bullet columns
    columns = [
        (Inches(0.55), "Challenges", case.get("challenges") or [], RED_ACCENT),
        (Inches(6.8), "The Joveo Solution", case.get("solution") or [], BLUE),
    ]
    col_top = Inches(3.15)
    for col_x, col_title, items, accent in columns:
        if not items:
            continue
        _add_textbox(
            slide,
            col_x,
            col_top,
            Inches(6.0),
            Inches(0.3),
            text=col_title,
            font_size=11,
            bold=True,
            color=NAVY,
        )
        _box, tf = _add_textbox(
            slide, col_x, col_top + Inches(0.34), Inches(6.0), Inches(1.9)
        )
        for item in items[:4]:
            p = tf.add_paragraph()
            p.space_before = Pt(1)
            p.space_after = Pt(3)
            rb = p.add_run()
            rb.text = "\u25cf  "
            _set_font(rb, size=8, color=accent)
            rt = p.add_run()
            rt.text = str(item)
            _set_font(rt, size=9, color=DARK_TEXT, name=FONT_BODY)

    # Bottom strip: numbered Next Steps 1-5, horizontal
    if next_steps:
        strip_top = Inches(5.6)
        _add_rounded_rect(
            slide, Inches(0.55), strip_top, Inches(12.25), Inches(1.25), LAVENDER_50
        )
        _add_textbox(
            slide,
            Inches(0.8),
            strip_top + Inches(0.08),
            Inches(3.0),
            Inches(0.28),
            text="Next Steps",
            font_size=11,
            bold=True,
            color=NAVY,
        )
        item_w = Inches(2.41)
        for idx, step in enumerate(next_steps[:5]):
            x = Inches(0.8) + idx * item_w
            chip = Inches(0.3)
            _add_oval(slide, x, strip_top + Inches(0.44), chip, chip, BLUE)
            _add_textbox(
                slide,
                x,
                strip_top + Inches(0.44),
                chip,
                chip,
                text=str(idx + 1),
                font_size=10,
                bold=True,
                color=WHITE,
                alignment=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE,
            )
            _box, _tf = _add_textbox(
                slide,
                x + Inches(0.38),
                strip_top + Inches(0.4),
                item_w - Inches(0.55),
                Inches(0.78),
                text=str(step),
                font_size=8,
                color=DARK_TEXT,
            )
            _set_body_font(_tf)

    _add_footer(slide, today)


# ===================================================================
# SLIDE N (Last) - Data Sources & Methodology
# ===================================================================


def _build_slide_data_sources(prs: Presentation, data: Dict):
    """Build a Data Sources slide listing all API sources, freshness, and confidence.

    This slide is appended as the LAST slide in the deck and provides
    transparency into the data pipeline that powered the media plan.
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    today = datetime.date.today().strftime("%B %d, %Y")

    # Off-white background
    _add_filled_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, OFF_WHITE)

    # Top band
    _add_top_band(slide, "DATA SOURCES & METHODOLOGY", today)

    # Action title
    _add_textbox(
        slide,
        Inches(0.55),
        Inches(0.92),
        Inches(12.2),
        Inches(0.5),
        text="Transparency into the data pipeline powering this media plan",
        font_size=15,
        bold=True,
        color=NAVY,
    )

    # S89: data-freshness line -- tells the reader exactly how current the
    # figures are ("Data current as of <date>").
    _add_textbox(
        slide,
        Inches(0.55),
        Inches(1.34),
        Inches(12.2),
        Inches(0.25),
        text=f"Data current as of {today}",
        font_size=9,
        italic=True,
        color=MUTED_TEXT,
    )

    # S89 KEYSTONE: "Joveo measured" callout -- only when the budget engine
    # attached first-party measured outcomes from the cg_benchmarks warehouse.
    # Read defensively; the key is absent in the common (no-match) case.
    _ba_meta = (data.get("_budget_allocation") or {})
    if isinstance(_ba_meta, dict):
        _ba_meta = _ba_meta.get("metadata") or {}
    else:
        _ba_meta = {}
    _real_outcomes = _ba_meta.get("real_outcomes") if isinstance(_ba_meta, dict) else None
    if _real_outcomes:
        try:
            _n_measured = len(_real_outcomes) if isinstance(_real_outcomes, list) else 1
        except Exception:  # noqa: BLE001
            _n_measured = 1
        _badge_w = Inches(3.0)
        _badge_x = SLIDE_WIDTH - _badge_w - Inches(0.55)
        _badge_top = Inches(0.95)
        _add_rounded_rect(slide, _badge_x, _badge_top, _badge_w, Inches(0.5), LAVENDER_100)
        _add_filled_rect(slide, _badge_x, _badge_top, Inches(0.06), Inches(0.5), GREEN)
        _badge_label = (
            f"Joveo measured · {_n_measured} role"
            f"{'s' if _n_measured != 1 else ''} matched"
        )
        _bbox, _btf = _add_textbox(
            slide,
            _badge_x + Inches(0.18),
            _badge_top,
            _badge_w - Inches(0.28),
            Inches(0.5),
            text=_badge_label,
            font_size=9,
            bold=True,
            color=NAVY,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    # Extract enrichment data
    enriched = data.get("_enriched", {})
    if not isinstance(enriched, dict):
        enriched = {}
    summary = enriched.get("enrichment_summary", {})
    if not isinstance(summary, dict):
        summary = {}

    apis_called = summary.get("apis_called") or []
    apis_succeeded = summary.get("apis_succeeded") or []
    apis_failed = summary.get("apis_failed") or []
    apis_skipped = summary.get("apis_skipped") or []
    api_details = summary.get("api_details", {})
    if not isinstance(api_details, dict):
        api_details = {}
    # Prefer the computed confidence from Excel's Sources sheet for consistency
    confidence_score = (
        data.get("_computed_confidence_pct") or summary.get("confidence_score") or 0
    )
    # Normalize: if stored as 0-100 integer, convert to 0-1 for _fmt_pct
    if isinstance(confidence_score, (int, float)) and confidence_score > 1:
        confidence_score = confidence_score / 100.0
    total_time = summary.get("total_time_seconds") or 0

    # ---- HERO STATS BAR ----
    bar_top = Inches(1.55)
    bar_h = Inches(0.85)
    _add_filled_rect(slide, Inches(0.55), bar_top, Inches(12.2), bar_h, NAVY)
    _add_filled_rect(slide, Inches(0.55), bar_top, Inches(12.2), Inches(0.03), TEAL)

    # Exclude skipped APIs from succeeded count
    truly_succeeded = set(apis_succeeded) - set(apis_skipped)
    hero_metrics = [
        (str(len(apis_called)), "APIs Called"),
        (str(len(truly_succeeded)), "APIs Succeeded"),
        (str(len(apis_failed)), "APIs Failed"),
    ]
    if apis_skipped:
        hero_metrics.append((str(len(apis_skipped)), "APIs Skipped"))
    hero_metrics.extend(
        [
            (
                _fmt_pct(confidence_score, decimals=0) if confidence_score else "—",
                "Confidence Score",
            ),
            (f"{total_time:.1f}s" if total_time else "—", "Fetch Time"),
        ]
    )

    hero_w = Inches(2.44)
    for i, (val, label) in enumerate(hero_metrics):
        hx = Inches(0.55) + i * hero_w
        _add_textbox(
            slide,
            hx,
            bar_top + Inches(0.05),
            hero_w,
            Inches(0.48),
            text=val,
            font_size=22,
            bold=True,
            color=TEAL,
            alignment=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        _add_textbox(
            slide,
            hx,
            bar_top + Inches(0.52),
            hero_w,
            Inches(0.25),
            text=label,
            font_size=8,
            bold=False,
            color=LIGHT_MUTED,
            alignment=PP_ALIGN.CENTER,
        )
        # Thin divider between metrics (skip after last)
        if i < len(hero_metrics) - 1:
            div_x = Inches(0.55) + (i + 1) * hero_w
            _add_filled_rect(
                slide,
                div_x,
                bar_top + Inches(0.15),
                Inches(0.015),
                Inches(0.55),
                RGBColor(0x1A, 0x45, 0x70),
            )

    # ---- DATA SOURCES TABLE ----
    table_top = Inches(2.65)
    table_left = Inches(0.55)
    table_w = Inches(12.2)
    row_h = Inches(0.35)

    # Section header
    _add_textbox(
        slide,
        table_left,
        table_top - Inches(0.4),
        Inches(4),
        Inches(0.35),
        text="API DATA SOURCES",
        font_size=11,
        bold=True,
        color=NAVY,
    )
    _add_filled_rect(
        slide, table_left, table_top - Inches(0.08), Inches(1.8), Inches(0.03), TEAL
    )

    # Table header row
    _add_filled_rect(slide, table_left, table_top, table_w, row_h, NAVY)
    col_widths = [Inches(2.8), Inches(2.5), Inches(2.0), Inches(2.2), Inches(2.7)]
    col_headers = ["Data Source", "Status", "Freshness", "Response Time", "Data Points"]
    col_offsets = [Inches(0)]
    for cw in col_widths[:-1]:
        col_offsets.append(col_offsets[-1] + cw)

    for ci, (header, cw, co) in enumerate(zip(col_headers, col_widths, col_offsets)):
        _add_textbox(
            slide,
            table_left + co + Inches(0.1),
            table_top,
            cw - Inches(0.1),
            row_h,
            text=header,
            font_size=9,
            bold=True,
            color=WHITE,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    # Build table rows from api_details
    # If api_details is empty, build from apis_called/succeeded/failed lists
    table_rows = []
    if api_details:
        for api_name, detail in api_details.items():
            if not isinstance(detail, dict):
                continue
            source = detail.get("source", "unknown")
            success = detail.get("success", False)
            elapsed = detail.get("elapsed_time") or 0
            status_label = detail.get("status", "unknown")

            # Determine freshness label
            if source == "live":
                freshness = "Live (real-time)"
            elif source == "cached":
                freshness = "Cached"
            elif source == "error":
                freshness = "--"
            else:
                freshness = "Curated"

            # Estimate data points from enriched data keys
            data_points = "--"
            result_key_map = {
                "BLS": "salary_data",
                "Adzuna": "job_market",
                "Census-ACS": "location_demographics",
                "BLS-QCEW": "industry_employment",
                "WorldBank": "global_indicators",
                "Clearbit": "clearbit_data",
                "Google-Ads": "google_ads_data",
                "Meta-Ads": "meta_ads_data",
                "Teleport": "teleport_data",
            }
            mapped_key = None
            for rk_prefix, rk_val in result_key_map.items():
                if rk_prefix.lower() in api_name.lower():
                    mapped_key = rk_val
                    break
            if mapped_key and enriched.get(mapped_key):
                raw = enriched[mapped_key]
                if isinstance(raw, dict):
                    data_points = str(len(raw))
                elif isinstance(raw, list):
                    data_points = str(len(raw))

            # Check if this API is in the skipped list
            is_skipped = api_name in apis_skipped

            # Determine status display
            # success=True but no actual data returned means it was effectively skipped
            has_data = (
                detail.get("data_points", 0) not in (0, None, "--")
                or data_points != "--"
            )
            if is_skipped or (success and not has_data):
                status_display = "Skipped"
                status_color = AMBER
            elif success:
                status_display = "Succeeded"
                status_color = GREEN
            elif status_label == "empty":
                status_display = "Skipped (no data)"
                status_color = AMBER
            elif status_label == "circuit_open":
                status_display = "Circuit Broken"
                status_color = RED_ACCENT
            else:
                status_display = "Failed"
                status_color = RED_ACCENT

            table_rows.append(
                (
                    api_name,
                    (status_display, status_color),
                    freshness,
                    f"{elapsed:.2f}s" if elapsed else "--",
                    data_points,
                )
            )
    else:
        # Fallback: build rows from simple lists
        # Check skipped BEFORE succeeded -- an API can appear in both
        all_apis = set(apis_called) if apis_called else set()
        for api_name in sorted(all_apis):
            if api_name in apis_skipped:
                status_display = "Skipped"
                status_color = AMBER
                freshness = "--"
            elif api_name in apis_succeeded:
                status_display = "Succeeded"
                status_color = GREEN
                freshness = "Live (real-time)"
            elif api_name in apis_failed:
                status_display = "Failed"
                status_color = RED_ACCENT
                freshness = "--"
            else:
                status_display = "Unknown"
                status_color = MUTED_TEXT
                freshness = "--"
            table_rows.append(
                (
                    api_name,
                    (status_display, status_color),
                    freshness,
                    "--",
                    "--",
                )
            )

    # Add SlotOps LinkedIn dataset as a data source when used
    _li_intel_ds = (data.get("_gold_standard") or {}).get("linkedin_intelligence", {})
    if not _li_intel_ds:
        _li_intel_ds = data.get("_slotops_linkedin_benchmarks", {})
    if _li_intel_ds and _li_intel_ds.get("sample_size", 0) > 0:
        _li_sample = _li_intel_ds.get(
            "total_jobs_analyzed", _li_intel_ds.get("sample_size", 108871)
        )
        table_rows.append(
            (
                "SlotOps LinkedIn (108K jobs)",
                ("Loaded", GREEN),
                "Curated (Apr 2025-Apr 2026)",
                "--",
                f"{_li_sample:,}" if isinstance(_li_sample, int) else str(_li_sample),
            )
        )

    # If no API data at all, show a placeholder row
    if not table_rows:
        table_rows.append(
            (
                "No API enrichment data",
                ("—", MUTED_TEXT),
                "Curated benchmarks used",
                "--",
                "--",
            )
        )

    # Render table rows (cap at 12 to fit the slide)
    max_display_rows = 12
    for ri, row_data in enumerate(table_rows[:max_display_rows]):
        ry = table_top + row_h * (ri + 1)
        bg = WHITE if ri % 2 == 0 else RGBColor(0xF8, 0xF6, 0xF3)
        _add_filled_rect(slide, table_left, ry, table_w, row_h, bg)
        # Thin ruled line
        _add_filled_rect(
            slide,
            table_left,
            ry + row_h - Inches(0.008),
            table_w,
            Inches(0.008),
            WARM_GRAY,
        )

        api_name_str, (status_str, status_clr), freshness_str, elapsed_str, dp_str = (
            row_data
        )

        vals = [api_name_str, status_str, freshness_str, elapsed_str, dp_str]
        colors = [DARK_TEXT, status_clr, MUTED_TEXT, MUTED_TEXT, DARK_TEXT]
        bolds = [True, False, False, False, False]

        for ci, (val, clr, bld, cw, co) in enumerate(
            zip(vals, colors, bolds, col_widths, col_offsets)
        ):
            _add_textbox(
                slide,
                table_left + co + Inches(0.1),
                ry,
                cw - Inches(0.1),
                row_h,
                text=str(val),
                font_size=8,
                bold=bld,
                color=clr,
                anchor=MSO_ANCHOR.MIDDLE,
            )

    # Overflow indicator
    if len(table_rows) > max_display_rows:
        overflow_y = table_top + row_h * (max_display_rows + 1) + Inches(0.05)
        _add_textbox(
            slide,
            table_left,
            overflow_y,
            table_w,
            Inches(0.2),
            text=f"+ {len(table_rows) - max_display_rows} additional sources (see appendix)",
            font_size=7,
            italic=True,
            color=MUTED_TEXT,
        )

    # ---- METHODOLOGY FOOTER ----
    # Thin rule line
    _add_rule_line(slide, 0.55, 7.0, 12.2, "D6CFC2")

    # S49 P2-20: Research-backed key recommendations (compact, bottom section)
    _research_recs = data.get("_research_recommendations") or []
    if _research_recs:
        _recs_text = "Key Insights: " + " | ".join(_research_recs[:5])
        _add_textbox(
            slide,
            Inches(0.55),
            Inches(6.65),
            Inches(12.2),
            Inches(0.35),
            text=_recs_text,
            font_size=7,
            bold=False,
            color=TEAL,
        )

    meth_text = (
        "Methodology: Data sourced from real-time API integrations, Nova AI Suite's proprietary "
        "job board knowledge base (91+ platforms), trend engine (4-year history), and "
        "curated industry benchmarks. Confidence score reflects the ratio of successful "
        "live API calls to total attempted."
    )
    _add_textbox(
        slide,
        Inches(0.55),
        Inches(7.05),
        Inches(10),
        Inches(0.35),
        text=meth_text,
        font_size=6,
        italic=True,
        color=LIGHT_MUTED,
    )

    # S89: data-source / provenance attribution line. Names the figure origins
    # so the deck is self-attesting; cites Joveo measured outcomes when present.
    _provenance_bits = [
        "Provenance: live APIs + Nova KB + curated benchmarks",
    ]
    if _real_outcomes:
        _provenance_bits.append("calibrated against Joveo measured outcomes (cg_benchmarks)")
    _provenance_text = " · ".join(_provenance_bits) + "."
    _add_textbox(
        slide,
        Inches(0.55),
        Inches(6.78),
        Inches(12.2),
        Inches(0.22),
        text=_provenance_text,
        font_size=6,
        italic=True,
        color=MUTED_TEXT,
    )

    # Generation timestamp
    _add_textbox(
        slide,
        Inches(10.5),
        Inches(7.05),
        Inches(2.2),
        Inches(0.2),
        text=f"Generated: {today}",
        font_size=7,
        bold=False,
        color=MUTED_TEXT,
        alignment=PP_ALIGN.RIGHT,
    )

    # Footer
    _add_footer(slide, today)


# ===================================================================
# CHART SLIDES - Matplotlib-generated visual data slides
# ===================================================================


def _build_slide_budget_pie_chart(prs: Presentation, data: Dict):
    """Build a slide with a matplotlib pie chart showing budget allocation by channel.

    This slide provides a visual breakdown of how the campaign budget is
    distributed across channels. Only added when matplotlib is available
    and channel data exists.
    """
    channels = _selected_channels(data)
    if not channels:
        return

    budget_alloc = data.get("_budget_allocation", {})
    ba_channel_alloc = (
        budget_alloc.get("channel_allocations", {})
        if isinstance(budget_alloc, dict)
        else {}
    )

    labels = []
    sizes = []

    # S6 fix: reconcile ALL channels' percentages together (largest-remainder
    # rounding) instead of rounding each independently, so this pie's legend
    # percentages sum to exactly 100 instead of drifting to 101%/99%.
    _reconciled_pct = _reconcile_channel_percentages(channels, ba_channel_alloc)

    for ch_key, ch_data in channels.items():
        label = ch_data.get("label", ch_key.replace("_", " ").title())
        pct = _reconciled_pct.get(ch_key, ch_data.get("pct") or 0)

        if pct > 0:
            labels.append(label)
            sizes.append(pct)

    if not labels:
        return

    chart_bytes = _generate_pie_chart_image(labels, sizes)
    if not chart_bytes:
        return

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    today = datetime.date.today().strftime("%B %d, %Y")

    # Background
    _add_filled_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, OFF_WHITE)
    _add_top_band(slide, "BUDGET ALLOCATION OVERVIEW", today)

    client = data.get("client_name", "Client")
    _add_textbox(
        slide,
        Inches(0.55),
        Inches(0.92),
        Inches(12.2),
        Inches(0.45),
        text=f"Visual breakdown of channel investment strategy for {client}",
        font_size=15,
        bold=True,
        color=NAVY,
    )

    # Insert the pie chart image
    chart_stream = io.BytesIO(chart_bytes)
    slide.shapes.add_picture(
        chart_stream,
        Inches(2.5),
        Inches(1.6),
        Inches(8.3),
        Inches(5.0),
    )

    _add_footer(slide, today)


def _build_slide_conversion_funnel(prs: Presentation, data: Dict):
    """Build a slide with a matplotlib funnel chart showing conversion metrics.

    Shows the recruitment funnel: Impressions -> Clicks -> Applications -> Hires.
    Data sourced from budget allocation projections and enriched analytics.
    """
    budget_alloc = data.get("_budget_allocation", {})
    if not isinstance(budget_alloc, dict):
        budget_alloc = {}
    ba_total_proj = budget_alloc.get("total_projected", {})
    if not isinstance(ba_total_proj, dict):
        ba_total_proj = {}
    ba_metadata = budget_alloc.get("metadata", {})
    if not isinstance(ba_metadata, dict):
        ba_metadata = {}
    ba_channel_alloc = budget_alloc.get("channel_allocations", {})
    if not isinstance(ba_channel_alloc, dict):
        ba_channel_alloc = {}

    # Gather funnel metrics -- S48: per-channel sum for consistency
    applications = int(ba_total_proj.get("applications") or 0)
    hires = sum(int(ch.get("projected_hires") or 0) for ch in ba_channel_alloc.values())
    if hires == 0:
        hires = int(ba_total_proj.get("hires") or 0)

    # Estimate impressions and clicks from budget if not provided
    total_budget = ba_metadata.get("total_budget") or 0
    if total_budget <= 0:
        budget_val = _parse_budget_number(data.get("budget") or "")
        total_budget = budget_val if budget_val else 0

    # Calculate clicks and impressions from channel data or defaults
    total_clicks = 0
    total_impressions = 0
    for _ch_key, ch_data in ba_channel_alloc.items():
        if isinstance(ch_data, dict):
            total_clicks += int(ch_data.get("projected_clicks") or 0)
            total_impressions += int(ch_data.get("projected_impressions") or 0)

    # Fallback: estimate from budget using industry averages
    if total_clicks <= 0 and total_budget > 0:
        avg_cpc = 1.50  # Industry average CPC
        total_clicks = int(total_budget / avg_cpc)
    if total_impressions <= 0 and total_clicks > 0:
        avg_ctr = 0.035  # Industry average CTR 3.5%
        total_impressions = int(total_clicks / avg_ctr)
    if applications <= 0 and total_clicks > 0:
        avg_apply_rate = 0.05  # Industry average apply rate 5%
        applications = int(total_clicks * avg_apply_rate)
    if hires <= 0 and applications > 0:
        avg_hire_rate = 0.08  # Industry average hire rate 8%
        hires = max(1, int(applications * avg_hire_rate))

    # Need at least impressions to show a funnel
    if total_impressions <= 0:
        return

    chart_bytes = _generate_funnel_chart_image(
        total_impressions, total_clicks, applications, hires
    )
    if not chart_bytes:
        return

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    today = datetime.date.today().strftime("%B %d, %Y")

    # Background
    _add_filled_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, OFF_WHITE)
    _add_top_band(slide, "RECRUITMENT CONVERSION FUNNEL", today)

    client = data.get("client_name", "Client")
    _add_textbox(
        slide,
        Inches(0.55),
        Inches(0.92),
        Inches(12.2),
        Inches(0.45),
        text=f"Projected candidate pipeline from awareness to hire for {client}",
        font_size=15,
        bold=True,
        color=NAVY,
    )

    # Insert the funnel chart image
    chart_stream = io.BytesIO(chart_bytes)
    slide.shapes.add_picture(
        chart_stream,
        Inches(2.2),
        Inches(1.6),
        Inches(8.8),
        Inches(5.0),
    )

    _add_footer(slide, today)


# ===================================================================
# Public API
# ===================================================================


def generate_pptx(data: Dict[str, Any]) -> bytes:
    """
    Generate a premium LinkedIn-inspired PowerPoint presentation.

    Args:
        data: Dictionary containing client information, industry details,
              channel selections, and campaign parameters. Expected keys:
              client_name, industry, budget, campaign_goals, target_roles/roles,
              work_environment, channel_categories, locations, experience_level.

    Returns:
        bytes: The .pptx file content as bytes, suitable for streaming
               to a client or writing to disk.

    Raises:
        ValueError: If required data fields are missing.
        RuntimeError: If presentation generation fails.
    """
    if data is None or not isinstance(data, dict):
        raise ValueError("Data must be a non-null dictionary.")

    # Ensure minimum required fields have sensible defaults
    data.setdefault("client_name", "Client")
    # Normalize client name casing (preserves known brands)
    data["client_name"] = _proper_client_name(data["client_name"] or "Client")
    data.setdefault("industry", "general_entry_level")
    data.setdefault("locations", [])
    # Frontend sends "target_roles" but PPT uses "roles" -- normalize
    if "target_roles" in data and "roles" not in data:
        data["roles"] = data["target_roles"]
    data.setdefault("roles", [])
    # Normalize roles: dicts -> strings (frontend may send [{"title": "..."}])
    if data["roles"] and isinstance(data["roles"][0], dict):
        data["roles"] = [
            r.get("title", r.get("role", str(r))) if isinstance(r, dict) else str(r)
            for r in data["roles"]
        ]
    data.setdefault("campaign_goals", [])
    data.setdefault("channel_categories", {})
    # Frontend sends "budget_range" but PPT reads "budget" -- normalize
    if data.get("budget_range") and not data.get("budget"):
        data["budget"] = data["budget_range"]
    data.setdefault("budget", "TBD")
    # Frontend sends work_environment as array -- normalize to string
    we = data.get("work_environment", "hybrid")
    if isinstance(we, list):
        data["work_environment"] = we[0] if we else "hybrid"
    data.setdefault("work_environment", "hybrid")

    # Null safety - replace None values with defaults
    for key, default in [
        ("client_name", "Client"),
        ("company_name", "Client"),
        ("industry", "general_entry_level"),
        ("budget", "TBD"),
        ("work_environment", "hybrid"),
    ]:
        if data.get(key) is None:
            data[key] = default
    # Ensure list fields are actual lists
    for key in ["locations", "roles", "target_roles", "campaign_goals", "competitors"]:
        val = data.get(key)
        if val is None:
            data[key] = []
        elif isinstance(val, str):
            data[key] = [val]
    # Ensure channel_categories is a dict
    cc = data.get("channel_categories")
    if cc is None:
        data["channel_categories"] = {}
    elif isinstance(cc, list):
        data["channel_categories"] = {
            (item.get("name") or "" if isinstance(item, dict) else str(item)): True
            for item in cc
        }

    # S49 P2-20: Inject research-backed recommendations into data for slides
    try:
        from research_constants import RESEARCH_FINDINGS_SHORT

        _existing_recs = data.get("_research_recommendations") or []
        if not _existing_recs:
            data["_research_recommendations"] = RESEARCH_FINDINGS_SHORT
    except ImportError:
        pass

    # Industry label mapping (single source of truth in shared_utils.py)
    if not data.get("industry_label"):
        data["industry_label"] = _SHARED_INDUSTRY_LABEL_MAP.get(
            data["industry"], data["industry"].replace("_", " ").title()
        )

    # S89: resolve the plan's currency once so every money figure renders in the
    # plan's own symbol (£/€/₹/...) instead of a hardcoded "$". Defaults to USD.
    _set_active_currency(data)

    try:
        prs = Presentation()

        # Set document metadata for GEO/SEO discoverability
        core_props = prs.core_properties
        client = data.get("client_name", "Client")
        industry_label = data.get(
            "industry_label", (data.get("industry") or "").replace("_", " ").title()
        )
        core_props.title = f"Recruitment Media Plan - {client}"
        core_props.author = "Nova AI Suite"
        core_props.subject = f"AI-generated recruitment advertising media plan for {client} in the {industry_label} industry"
        core_props.keywords = f"recruitment media plan, {industry_label}, job advertising, programmatic recruitment, {client}, talent acquisition, hiring strategy"
        core_props.comments = f"Generated by Nova AI Media Plan Generator (media-plan-generator.onrender.com). Data sourced from 25 real-time APIs, 91+ job board platforms, and Nova AI Suite industry knowledge base."
        core_props.category = "Recruitment Advertising"
        core_props.last_modified_by = "Nova AI Suite"
        # craft:both#2: stamp the REAL generation timestamp -- python-pptx's
        # default template otherwise leaves docProps/core.xml's
        # created/modified at its own hardcoded 2013-01-27 template date,
        # byte-identical across every bundle regardless of when it was
        # actually generated (the xlsx sibling already gets this right).
        _now_utc = datetime.datetime.now(datetime.timezone.utc)
        core_props.created = _now_utc
        core_props.modified = _now_utc

        # Set 16:9 widescreen dimensions
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT

        # S48 note superseded: ~11-slide deck mirroring the official Joveo
        # client deck flow. Extended analysis still lives in Excel.
        # Slide order: Cover, Exec Summary, [Methodology], [Push Meets Pull],
        #              Channel Strategy, Budget+Pie (or Quality), [CPA Reference],
        #              Competitive Landscape, Timeline, Risk,
        #              [Case Study & Next Steps].
        # Bracketed narrative slides come from the deck KB and self-skip
        # when data/joveo_media_plan_deck_2026.json is unavailable.

        # Deck KB: narrative content shared with the official client deck
        deck = _load_deck_kb()

        # Slide 1: Premium cover / section divider
        _build_slide_cover(prs, data)

        # Quality warning disclaimer on cover slide (when enrichment degraded)
        # S49 Issue 17: Severity-scaled prefix -- "Caution" for minimal data,
        # "Warning" for limited data, "Note" for moderate degradation.
        # S1 (2026-07-03): this is an internal QA signal (enrichment confidence),
        # not a client-facing disclaimer -- gate it behind internal_qc_mode so it
        # never ships on a client cover slide by default.
        _quality_warn = data.get("quality_warning") or ""
        if _quality_warn and prs.slides and _internal_qc_mode(data):
            try:
                _cover_slide = prs.slides[0]
                # Determine severity prefix and color from warning text
                if "Minimal data" in _quality_warn:
                    _warn_prefix = "Caution"
                    _warn_color = "FF4444"  # Red for severe
                elif "Limited data" in _quality_warn:
                    _warn_prefix = "Warning"
                    _warn_color = "FF8C00"  # Dark orange for moderate
                else:
                    _warn_prefix = "Note"
                    _warn_color = "B7669E"  # Magenta for mild
                _add_textbox(
                    _cover_slide,
                    Inches(0.6),
                    Inches(7.0),
                    Inches(10),
                    Inches(0.3),
                    text=f"{_warn_prefix}: {_quality_warn}",
                    font_size=7,
                    italic=True,
                    color=_warn_color,
                )
            except Exception as _qw_err:
                logger.debug(
                    "Quality warning disclaimer failed (non-fatal): %s", _qw_err
                )

        # Slide 2: Executive Summary (hero stats + SCR + market context)
        _build_slide_executive_summary(prs, data)

        # Slide 3: Our Methodology (deck narrative -- never blocks generation)
        try:
            _build_slide_methodology(prs, data, deck)
        except Exception as _meth_exc:
            logger.debug("Methodology slide failed (non-fatal): %s", _meth_exc)

        # Slide 4: Push Meets Pull (deck narrative -- never blocks generation)
        try:
            _build_slide_push_meets_pull(prs, data, deck)
        except Exception as _pmp_exc:
            logger.debug("Push Meets Pull slide failed (non-fatal): %s", _pmp_exc)

        # Slide 5: Channel Strategy with benchmarks + attribution
        _build_slide_channel_strategy(prs, data)

        # Slide 6: Budget Allocation & ROI with embedded pie chart
        #          (or Quality Outcomes as fallback when budget_allocation is empty)
        budget_alloc_data = data.get("_budget_allocation", {})
        if isinstance(budget_alloc_data, dict) and budget_alloc_data:
            _ba_has_data = (
                (budget_alloc_data.get("metadata", {}).get("total_budget") or 0) > 0
                or budget_alloc_data.get("total_projected", {})
                or budget_alloc_data.get("channel_allocations", {})
            )
            if _ba_has_data:
                _build_slide_budget_allocation(prs, data)
                # NOTE: the embedded pie chart was removed -- it was dropped at a
                # fixed bottom-right position that overlapped the budget table, the
                # ROI callout, and the footer. Allocation is already shown by the
                # table's Allocation% column and the channel-mix bars on the
                # Channel Strategy slide, so a clean full-width table + dark
                # takeaway band (reference-deck style) reads better.
            else:
                _build_slide_quality_outcomes(prs, data)
        else:
            _build_slide_quality_outcomes(prs, data)

        # Slide 7: CPA Reference — industry-keyed KB content (only AI-training
        # is populated); _build_slide_cpa_reference itself drops the slide
        # entirely when this plan's industry key has no KB content.
        _slides_before_cpa = len(prs.slides)
        try:
            _build_slide_cpa_reference(prs, data, deck)
        except Exception as _cpa_exc:
            logger.debug("CPA Reference slide failed (non-fatal): %s", _cpa_exc)

        # visual:atria#4: when the CPA Reference slot was dropped (no KB
        # content for this plan's industry) and the plan names >=4 roles
        # with gold-standard role data, fill the empty slot with a compact
        # role-breakdown table instead of shipping a multi-role plan with
        # no role-level view anywhere in the deck.
        if len(prs.slides) == _slides_before_cpa:
            try:
                _build_slide_role_breakdown(prs, data)
            except Exception as _rb_exc:
                logger.debug("Role breakdown slide failed (non-fatal): %s", _rb_exc)

        # Slide 8: Competitive Landscape (always shown)
        _build_slide_competitive_landscape(prs, data)

        # Slide 9: Implementation Timeline + Competitive Context
        _build_slide_comparison_timeline(prs, data)

        # Slide 10: Risk Analysis -- always included for C-suite readiness
        _build_slide_risk_analysis(prs, data)

        # Slide 11: The Joveo Advantage -- value-prop close (non-fatal)
        try:
            _build_slide_why_joveo(prs, data, deck)
        except Exception as _wj_exc:
            logger.debug("The Joveo Advantage slide failed (non-fatal): %s", _wj_exc)

        # Slide 12: Case Study & Next Steps (deck narrative -- never blocks generation)
        try:
            _build_slide_case_study_next_steps(prs, data, deck)
        except Exception as _cs_exc:
            logger.debug(
                "Case Study & Next Steps slide failed (non-fatal): %s", _cs_exc
            )

        # Data Sources & Methodology -- REMOVED from client deck (S50)
        # _build_slide_data_sources(prs, data)  # kept for internal debugging

        _n_slides_final = len(prs.slides)
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        # Embed the brand font so the deck renders in Poppins everywhere (not a
        # serif/Calibri substitute on machines lacking the font).
        _pptx_bytes = _embed_fonts_in_pptx(buffer.getvalue())
        # craft:both#3 / craft:both#6: theme default typeface + app.xml
        # slide count/signature.
        return _fix_pptx_package_hygiene(_pptx_bytes, _n_slides_final)

    except Exception as exc:
        logger.error("generate_pptx top-level crash: %s", exc, exc_info=True)
        # Return a minimal error presentation so the caller always gets valid bytes
        try:
            err_prs = Presentation()
            err_prs.slide_width = SLIDE_WIDTH
            err_prs.slide_height = SLIDE_HEIGHT
            err_slide = err_prs.slides.add_slide(err_prs.slide_layouts[6])
            _add_filled_rect(
                err_slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, NAVY
            )
            _add_textbox(
                err_slide,
                Inches(1),
                Inches(2),
                Inches(11),
                Inches(1),
                text="Media Plan Generation Error",
                font_size=28,
                bold=True,
                color="FFFFFF",
            )
            _add_textbox(
                err_slide,
                Inches(1),
                Inches(3.5),
                Inches(11),
                Inches(1.5),
                text=f"An error occurred while generating the presentation: {exc}\n\n"
                "Please try again or contact support if the issue persists.",
                font_size=14,
                color="CCCCCC",
            )
            err_buf = io.BytesIO()
            err_prs.save(err_buf)
            err_buf.seek(0)
            return err_buf.getvalue()
        except Exception as inner_exc:
            logger.error(
                "generate_pptx: even error presentation creation failed: %s",
                inner_exc,
                exc_info=True,
            )
            raise RuntimeError(
                f"Failed to generate PowerPoint presentation: {exc}"
            ) from exc


# ===================================================================
# CLI entry point for testing
# ===================================================================

if __name__ == "__main__":
    sample_data = {
        "client_name": "Acme Healthcare",
        "client_website": "https://www.acmehealthcare.com",
        "industry": "healthcare_medical",
        "industry_label": "Healthcare & Medical",
        "locations": [
            "New York, NY",
            "Chicago, IL",
            "Houston, TX",
            "Phoenix, AZ",
            "San Diego, CA",
        ],
        "roles": [
            "Registered Nurse",
            "Medical Assistant",
            "Physical Therapist",
            "Lab Technician",
            "Pharmacist",
        ],
        "job_categories": ["Clinical", "Allied Health", "Administrative"],
        "use_case": "Scaling clinical hiring across 5 metro areas to meet Q3 demand surge",
        "campaign_goals": ["high_volume", "cost_efficiency", "speed_to_hire"],
        "work_environment": "on_site",
        "budget": "$75,000 / month",
        "competitors": ["HCA Healthcare", "UnitedHealth Group"],
        "channel_categories": {
            "regional_boards": True,
            "global_boards": True,
            "niche_boards": True,
            "social_media": True,
            "programmatic_dsp": True,
            "employer_branding": False,
            "apac_regional": False,
            "emea_regional": False,
        },
        "include_dei": True,
        "include_innovative": False,
        "include_budget_guide": True,
        "include_global_supply": False,
    }

    pptx_bytes = generate_pptx(sample_data)
    output_path = "media_plan_sample.pptx"
    with open(output_path, "wb") as f:
        f.write(pptx_bytes)
    print(f"Generated {output_path} ({len(pptx_bytes):,} bytes)")
