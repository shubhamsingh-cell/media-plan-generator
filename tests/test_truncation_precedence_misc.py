"""Repo-wide sweep (follow-up to 8fc9a4f): `X or ""[:N]` / `X or [][:N]`
operator-precedence no-ops across the remaining modules, verified by targeted
repro before fixing.

Python precedence: slicing binds tighter than `or`, so the slice applies to
the empty LITERAL and the intended cap is silently a no-op. Sites fixed and
covered here:

  1. api_portal._get_dashboard_summary day key
     (`entry.get("timestamp") or ""[:10]`) -- the [:10] date-prefix cut never
     applied, so each full ISO timestamp became its own "day" bucket and
     recent_daily_usage degenerated into one row per request. Now
     `(entry.get("timestamp") or "")[:10]`. ADJACENT FIX in the same
     function, found while verifying this consumer: the line above it,
     `e.get("timestamp") or "" >= cutoff`, parsed as
     `ts or ("" >= cutoff)` -- comparison binds tighter than `or` -- so the
     "last 7 days" filter passed EVERY entry that had any timestamp at all.
     Now `(e.get("timestamp") or "") >= cutoff`.
  2. api_enrichment.fetch_onet_occupation_data technology skills
     (`for ex in cat.get("example") or [][:2]`) -- the 2-examples-per-category
     cap was a no-op (repro: 6 examples per category all flowed through until
     the overall [:10] backstop). Now `(cat.get("example") or [])[:2]`.
  3. llm_router `__main__` CLI demo print (`result.get('text') or ''[:500]`)
     -- debug-only output on an unreachable-from-import branch; verified by
     reading and fixed to `(result.get('text') or '')[:500]`. No behavioral
     test (no import-reachable seam); the source-level guard in
     tests/test_truncation_precedence_fixes.py covers regression.
  4. quick_plan.get_role_insights preferred_platforms
     (`strategy.get("preferred_platforms") or [][:5]`) -- blue-collar's
     8-platform list flowed through uncapped. Now capped at 5.
  5. social_plan.generate_creative_briefs recommended_formats
     (`platform.get("ad_formats") or [][:4]`) -- Facebook's 7 formats all
     flowed to the briefs UI tag list. Now capped at 4.
  6. social_plan.generate_social_plan_excel Content Calendar rows
     (`for post in week.get("posts") or [][:20]`) -- the per-week 20-row
     limit promised by the inline comment never applied. Now capped at 20.
  7. ppt_generator._build_slide_geopolitical_risk event text
     (`ev.get("event") or ""[:80]`) -- RENDERED SLIDE TEXT; a 169-char event
     rendered in full (repro before fix). Fixed with the module's own
     word-boundary helper `_trunc_word(..., 80)` -- NOT a hard slice --
     matching how 8fc9a4f handled rendered text in excel_v2.
  8. archive/excel_legacy.py Joveo DSP / Social Media publisher lists
     (two `... or [][:10]` sites) -- module is live via the app.py legacy
     import path; both parenthesized. Covered by the import smoke test here
     plus the widened source-level guard (no behavioral repro required).

Runs under pytest, or standalone:
``python3 tests/test_truncation_precedence_misc.py``.
"""

from __future__ import annotations

import datetime as _dt
import io
import json
import os
import sys
import urllib.request
from pathlib import Path

PROJECT_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(PROJECT_ROOT))

import openpyxl  # noqa: E402
from pptx import Presentation  # noqa: E402

import api_enrichment  # noqa: E402
import api_portal  # noqa: E402
import collar_intelligence  # noqa: E402
import ppt_generator as ppt  # noqa: E402
import quick_plan  # noqa: E402
import social_plan  # noqa: E402


# ---------------------------------------------------------------------------
# Shared helpers
# ---------------------------------------------------------------------------
def _assert_trunc_word_boundary_cut(cell: str, source: str, cap: int):
    """`cell` must be a clean word-boundary prefix of `source` produced by
    ppt_generator._trunc_word: ends in exactly '...', never longer than
    cap + 3, and the cut never lands mid-word."""
    assert cell.endswith("..."), f"expected trailing ellipsis: {cell!r}"
    assert not cell.endswith("...."), f"doubled ellipsis: {cell!r}"
    assert len(cell) <= cap + 3, f"{len(cell)} chars > cap {cap} + ellipsis"
    body = cell[:-3]
    assert body, "empty body before ellipsis"
    assert source.startswith(body), f"not a prefix of the source: {body!r}"
    # The next source character after the kept body must not be mid-word.
    assert not source[len(body)].isalnum(), f"mid-word cut: ...{body[-15:]!r}"


def _sheet_strings(ws):
    for row in ws.iter_rows(values_only=True):
        for val in row:
            if isinstance(val, str):
                yield val


# ===========================================================================
# Site 1: api_portal dashboard -- day-key grouping + last-7-days filter
# ===========================================================================
def _dashboard_summary_with_log(usage_log):
    original = api_portal._load_api_keys
    api_portal._load_api_keys = lambda: {"keys": {}, "usage_log": usage_log}
    try:
        return api_portal._get_dashboard_summary()
    finally:
        api_portal._load_api_keys = original


def test_dashboard_daily_usage_groups_by_date_prefix():
    """Two same-day entries must share ONE date bucket keyed by the 10-char
    ISO date prefix -- pre-fix, each full timestamp was its own 'day'."""
    day = (_dt.datetime.utcnow() - _dt.timedelta(days=1)).strftime("%Y-%m-%d")
    summary = _dashboard_summary_with_log(
        [
            {"timestamp": f"{day}T10:07:48.123456Z"},
            {"timestamp": f"{day}T18:30:00.000000Z"},
        ]
    )
    assert summary["recent_daily_usage"] == [{"date": day, "requests": 2}]


def test_dashboard_recent_usage_honours_seven_day_cutoff():
    """Adjacent fix on the line above the day key: the cutoff comparison used
    to be truthiness-of-timestamp, so 30-day-old entries leaked into the
    'Last 7 days' series. Missing timestamps stay excluded."""
    now = _dt.datetime.utcnow()
    recent_day = (now - _dt.timedelta(days=1)).strftime("%Y-%m-%d")
    old_ts = (now - _dt.timedelta(days=30)).isoformat() + "Z"
    summary = _dashboard_summary_with_log(
        [
            {"timestamp": f"{recent_day}T09:00:00.000000Z"},
            {"timestamp": old_ts},
            {"endpoint": "/api/x"},  # no timestamp at all -- must not crash
        ]
    )
    assert summary["recent_daily_usage"] == [{"date": recent_day, "requests": 1}]


def test_dashboard_empty_usage_log_safe():
    summary = _dashboard_summary_with_log([])
    assert summary["recent_daily_usage"] == []


# ===========================================================================
# Site 2: api_enrichment O*NET technology skills -- 2 examples per category
# ===========================================================================
class _FakeOnetResponse:
    def __init__(self, payload):
        self._payload = json.dumps(payload).encode()

    def read(self):
        return self._payload

    def __enter__(self):
        return self

    def __exit__(self, *args):
        return False


def _fake_onet_urlopen(req, timeout=10, context=None):
    url = req.full_url
    if url.endswith("/summary/technology_skills"):
        return _FakeOnetResponse(
            {
                "category": [
                    {"title": "Cat A", "example": [{"name": f"ToolA{i}"} for i in range(6)]},
                    {"title": "Cat B", "example": [{"name": f"ToolB{i}"} for i in range(6)]},
                    {"title": "Cat C"},  # no "example" key -- must be safe
                ]
            }
        )
    if url.endswith("/summary/skills") or url.endswith("/summary/knowledge"):
        return _FakeOnetResponse({"element": []})
    return _FakeOnetResponse({"title": "Software Engineer", "description": "d"})


def test_onet_technology_skills_capped_at_two_per_category():
    """Pre-fix repro: 6 examples per category all flowed through (12 total,
    then the [:10] backstop). Post-fix: 2 per category."""
    saved = (
        api_enrichment._get_cached,
        api_enrichment._set_cached,
        urllib.request.urlopen,
        os.environ.get("ONET_API_KEY"),
    )
    api_enrichment._get_cached = lambda key: None
    api_enrichment._set_cached = lambda key, data: None
    urllib.request.urlopen = _fake_onet_urlopen
    os.environ["ONET_API_KEY"] = "fake-key-for-test"
    try:
        result = api_enrichment.fetch_onet_occupation_data(["software engineer"])
    finally:
        api_enrichment._get_cached = saved[0]
        api_enrichment._set_cached = saved[1]
        urllib.request.urlopen = saved[2]
        if saved[3] is None:
            os.environ.pop("ONET_API_KEY", None)
        else:
            os.environ["ONET_API_KEY"] = saved[3]

    occ = result["occupations"]["software engineer"]
    assert occ["technology_skills"] == ["ToolA0", "ToolA1", "ToolB0", "ToolB1"]


# ===========================================================================
# Site 4: quick_plan preferred_platforms -- capped at 5
# ===========================================================================
def test_role_insights_preferred_platforms_capped_at_five():
    insights = quick_plan.get_role_insights("Forklift Operator", "logistics_supply_chain")
    assert insights["collar_type"] == "blue_collar"
    full = collar_intelligence.COLLAR_STRATEGY["blue_collar"]["preferred_platforms"]
    assert len(full) > 5, "fixture must exceed the cap to prove it applies"
    assert insights["preferred_platforms"] == full[:5]


def test_role_insights_preferred_platforms_short_list_unchanged():
    insights = quick_plan.get_role_insights("Registered Nurse", "healthcare_medical")
    assert insights["collar_type"] == "grey_collar"
    full = collar_intelligence.COLLAR_STRATEGY["grey_collar"]["preferred_platforms"]
    assert len(full) <= 5
    assert insights["preferred_platforms"] == full


def test_role_insights_missing_strategy_yields_empty_platforms():
    class _StubCollarIntel:
        COLLAR_STRATEGY = {}

        @staticmethod
        def classify_collar(role, industry):
            return {"collar_type": "unclassified", "confidence": 0.4}

    saved = quick_plan._collar_intel
    quick_plan._collar_intel = _StubCollarIntel()
    try:
        insights = quick_plan.get_role_insights("Mystery Role", "general_entry_level")
    finally:
        quick_plan._collar_intel = saved
    assert insights["preferred_platforms"] == []


# ===========================================================================
# Site 5: social_plan creative-brief recommended_formats -- capped at 4
# ===========================================================================
def test_creative_brief_formats_capped_at_four():
    briefs = social_plan.generate_creative_briefs(
        ["facebook"], "Warehouse Associate", "logistics_supply_chain", "blue_collar"
    )
    full = social_plan.SOCIAL_PLATFORMS["facebook"]["ad_formats"]
    assert len(full) > 4, "fixture must exceed the cap to prove it applies"
    assert briefs[0]["recommended_formats"] == full[:4]


def test_creative_brief_formats_short_list_unchanged():
    briefs = social_plan.generate_creative_briefs(
        ["indeed_sponsored"], "Warehouse Associate", "logistics_supply_chain", "blue_collar"
    )
    full = social_plan.SEARCH_PLATFORMS["indeed_sponsored"]["ad_formats"]
    assert len(full) <= 4
    assert briefs[0]["recommended_formats"] == full


def test_creative_brief_missing_ad_formats_safe():
    social_plan.SOCIAL_PLATFORMS["_test_stub"] = {"name": "Stub Platform"}
    try:
        briefs = social_plan.generate_creative_briefs(
            ["_test_stub"], "Warehouse Associate", "logistics_supply_chain", "blue_collar"
        )
    finally:
        del social_plan.SOCIAL_PLATFORMS["_test_stub"]
    assert briefs[0]["recommended_formats"] == []


# ===========================================================================
# Site 6: social_plan Excel Content Calendar -- 20 rows per week
# ===========================================================================
_BASE_SOCIAL_PLAN = None


def _base_social_plan():
    global _BASE_SOCIAL_PLAN
    if _BASE_SOCIAL_PLAN is None:
        _BASE_SOCIAL_PLAN = social_plan.generate_social_media_plan(
            role="Warehouse Associate",
            location="Dallas, TX",
            industry="logistics_supply_chain",
            budget=200000,
            duration_weeks=1,
        )
    return _BASE_SOCIAL_PLAN


def _calendar_week(week_num, theme, n_posts):
    return {
        "week": week_num,
        "theme": theme,
        "posts": [
            {
                "platform": "facebook",
                "platform_name": "Facebook",
                "content_type": f"Post {i}",
                "format": "Image",
                "day": "Monday",
                "expected_engagement": "medium",
            }
            for i in range(n_posts)
        ],
    }


def _calendar_row_counts(xlsx_bytes):
    """Count rendered data rows under each 'Week N:' header, stopping at the
    first blank row so the footer never bleeds into the last week."""
    wb = openpyxl.load_workbook(io.BytesIO(xlsx_bytes))
    ws = wb["Content Calendar"]
    counts = {}
    current = None
    for row in ws.iter_rows(values_only=True):
        b_val = row[1] if len(row) > 1 else None
        if isinstance(b_val, str) and b_val.startswith("Week "):
            current = b_val
            counts[current] = 0
            continue
        if b_val in (None, ""):
            current = None
            continue
        if isinstance(b_val, str) and b_val == "Platform":
            continue  # column-header row
        if current is not None:
            counts[current] += 1
    return counts


def test_content_calendar_excel_caps_posts_at_twenty_per_week():
    plan = dict(_base_social_plan())
    plan["content_calendar"] = [
        _calendar_week(1, "Over Cap", 25),
        _calendar_week(2, "Under Cap", 3),
    ]
    counts = _calendar_row_counts(social_plan.generate_social_plan_excel(plan))
    assert counts["Week 1: Over Cap"] == 20, "the 20-post cap is still a no-op"
    assert counts["Week 2: Under Cap"] == 3, "short week must render unchanged"


def test_content_calendar_excel_empty_posts_safe():
    plan = dict(_base_social_plan())
    plan["content_calendar"] = [{"week": 1, "theme": "No Posts"}]  # no "posts" key
    counts = _calendar_row_counts(social_plan.generate_social_plan_excel(plan))
    assert counts["Week 1: No Posts"] == 0


# ===========================================================================
# Site 7: ppt_generator geopolitical event text -- word-boundary cap at 80
# ===========================================================================
_LONG_EVENT = (
    "Port congestion continues to disrupt inbound freight schedules across "
    "the region, creating downstream delays for warehouse staffing plans and "
    "last-mile scheduling teams."
)  # 169 chars -- pre-fix repro rendered this in full on the slide


def _geo_slide_texts(events):
    data = {
        "client_name": "Acme Corp",
        "_synthesized": {
            "geopolitical_context": {
                "risk_level": "moderate",
                "overall_risk_score": 5.0,
                "summary": "Elevated risk this quarter.",
                "recommendations": [],
                "locations": {
                    "Dallas, TX": {
                        "risk_score": 5.0,
                        "budget_adjustment_factor": 1.0,
                        "events": events,
                    }
                },
            }
        },
    }
    prs = Presentation()
    prs.slide_width = ppt.SLIDE_WIDTH
    prs.slide_height = ppt.SLIDE_HEIGHT
    ppt._build_slide_geopolitical_risk(prs, data)
    assert len(prs.slides) == 1
    return [
        s.text_frame.text
        for s in prs.slides[0].shapes
        if s.has_text_frame and s.text_frame.text.strip()
    ]


def test_geo_risk_event_text_capped_at_word_boundary():
    texts = _geo_slide_texts([{"event": _LONG_EVENT, "severity": "moderate"}])
    hits = [t for t in texts if "Port congestion" in t]
    assert hits, "event card did not render"
    rendered = hits[0]
    assert _LONG_EVENT not in rendered, "the 80-char cap is still a no-op"
    # Slice off the severity-icon prefix before asserting on the event text.
    idx = rendered.find("Port congestion")
    _assert_trunc_word_boundary_cut(rendered[idx:], _LONG_EVENT, 80)


def test_geo_risk_short_event_text_unchanged():
    short = "Minor road closures downtown."
    texts = _geo_slide_texts([{"event": short, "severity": "low"}])
    hits = [t for t in texts if short in t]
    assert hits, "event card did not render"
    assert hits[0].endswith(short), "short event must render in full, no ellipsis"


def test_geo_risk_missing_event_key_safe():
    texts = _geo_slide_texts([{"severity": "low"}])
    # The card still builds past the events loop: location name renders.
    assert any("Dallas, TX" in t for t in texts)


# ===========================================================================
# Site 8: archive/excel_legacy -- live legacy module must stay importable
# ===========================================================================
def test_archive_excel_legacy_imports():
    import importlib

    mod = importlib.import_module("archive.excel_legacy")
    assert hasattr(mod, "generate_excel")


if __name__ == "__main__":
    _failures = 0
    for _name, _fn in sorted(globals().items()):
        if _name.startswith("test_") and callable(_fn):
            try:
                _fn()
                print(f"PASS {_name}")
            except AssertionError as exc:
                _failures += 1
                print(f"FAIL {_name}: {exc}")
            except Exception as exc:  # noqa: BLE001
                _failures += 1
                print(f"ERROR {_name}: {exc}")
    if _failures:
        sys.exit(1)
