"""Tests for Agent W3B's ppt_generator.py bundle-quality fixes (rescore2).

Covers:
    1. visual:manpower#2 / consistency:manpower#1 / consistency:atria#1:
       the Push Meets Pull card itemizes EVERY push/pull channel, and the
       itemized dollar figures foot exactly to the printed total.
    2. consistency:atria#3 / consistency:manpower#4: the deck-only
       "Reach Multiplier" / "Channel Diversity Score" metrics (and the
       unsourced "18% churn" claim) are gone from the comparison slide /
       complications list.
    3. consistency:atria#2: the Role Breakdown slide's "Est. Median
       Salary" column reads the SAME ``_gold_standard`` per-role salary
       data the workbook renders, not ``_enriched.salary_data``.
    4. copy:both#4: the Next Steps slide interpolates client name, roles,
       locations, budget, and duration instead of being byte-identical
       across clients.
    5. copy:both#1/#2, visual:atria#2: competitor cards are tagged
       "Industry competitor" vs "Talent-market competitor" and the Why
       line varies by that classification (Amazon vs. Brookdale case).

Runs under pytest, or standalone: ``python3 tests/test_ppt_w3b_bundle_quality.py``.
"""

from __future__ import annotations

import io
import os
import re
import sys

import pytest

PROJECT_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, PROJECT_ROOT)

import ppt_generator as ppt  # noqa: E402
import display_format as fmt  # noqa: E402
from pptx import Presentation  # noqa: E402


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def _new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = ppt.SLIDE_WIDTH
    prs.slide_height = ppt.SLIDE_HEIGHT
    return prs


def _all_slide_text(prs: Presentation) -> list[str]:
    out: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                out.append(shape.text_frame.text)
    return out


def _money_tokens(text: str) -> list[float]:
    """Parse every '$X.XK'/'$X'/'$X,XXX' token in ``text`` into a float."""
    out = []
    for tok in re.findall(r"\$[\d,]+\.?\d*[KM]?", text):
        body = tok.lstrip("$").replace(",", "")
        mult = 1.0
        if body.endswith("K"):
            mult = 1_000.0
            body = body[:-1]
        elif body.endswith("M"):
            mult = 1_000_000.0
            body = body[:-1]
        out.append(float(body) * mult)
    return out


def _manpower_channel_allocations() -> dict:
    # Amounts mirror the real bundle: 4 large push channels + Social Media
    # (small push channel previously dropped by the [:4] cap) + Employer
    # Branding (pull).
    return {
        "programmatic_dsp": {"dollar_amount": 39535.5},
        "regional_boards": {"dollar_amount": 36175.5},
        "global_boards": {"dollar_amount": 33655.5},
        "niche_boards": {"dollar_amount": 29221.5},
        "social_media": {"dollar_amount": 3912.0},
        "employer_branding": {"dollar_amount": 7500.0},
    }


def _logistics_plan() -> dict:
    roles = [{"title": "CDL A Driver", "count": 100, "tier": "Hourly"}]
    return {
        "client_name": "Manpower - Amerigas",
        "industry": "logistics_supply_chain",
        "industry_label": "Logistics & Supply Chain",
        "budget": "$150,000",
        "campaign_duration": "6 months",
        "locations": ["Massachusetts", "Maine", "New Hampshire", "Rhode Island", "Connecticut", "Denver, CO"],
        "roles": ["CDL A Driver"],
        "target_roles": roles,
        "work_environment": "on_site",
        "hire_volume": "500 hires",
        "competitors": ["FedEx", "UPS", "XPO Logistics"],
        "_budget_allocation": {
            "channel_allocations": _manpower_channel_allocations(),
            "metadata": {"total_budget": 150_000},
            "total_projected": {"hires": 48, "applications": 18950},
        },
    }


def _healthcare_plan() -> dict:
    roles = [
        {"title": "Memory Care Associate", "count": 30, "tier": "Hourly"},
        {"title": "Nurse", "count": 20, "tier": "Professional"},
        {"title": "Cook", "count": 10, "tier": "Hourly"},
        {"title": "Driver", "count": 10, "tier": "Hourly"},
    ]
    return {
        "client_name": "atria Senior living",
        "industry": "healthcare_medical",
        "industry_label": "Healthcare & Medical",
        "budget": "$300,000",
        "campaign_duration": "18 months",
        "locations": ["New York, NY"],
        "roles": [r["title"] for r in roles],
        "target_roles": roles,
        "work_environment": "on_site",
        "hire_volume": "500+ hires",
        "competitors": ["Brookdale Senior Living", "Sunrise Senior Living", "Amazon"],
        "_budget_allocation": {
            "channel_allocations": {
                "programmatic_dsp": {"dollar_amount": 84_336.0},
                "global_boards": {"dollar_amount": 79_400.0},
                "regional_boards": {"dollar_amount": 56_400.0},
                "niche_boards": {"dollar_amount": 35_000.0},
                "social_media": {"dollar_amount": 20_850.0},
                "employer_branding": {"dollar_amount": 24_000.0},
            },
            "metadata": {"total_budget": 300_000},
            "total_projected": {"hires": 57, "applications": 17780},
        },
    }


# ---------------------------------------------------------------------------
# 1. Push Meets Pull footing
# ---------------------------------------------------------------------------
class TestPushMeetsPullFooting:
    def test_split_line_itemizes_every_channel(self):
        data = _logistics_plan()
        push, pull = ppt._push_pull_channel_split(data)
        # All 5 push channels (incl. the small Social Media one) present.
        push_labels = {label for label, _ in push}
        assert len(push) == 5
        assert any("Social" in label for label in push_labels)

    def test_split_line_itemized_sum_equals_printed_total(self):
        data = _logistics_plan()
        push, pull = ppt._push_pull_channel_split(data)
        total_budget = sum(d for _, d in push) + sum(d for _, d in pull)
        line = ppt._push_pull_split_line(push, total_budget)

        tokens = _money_tokens(line)
        assert len(tokens) >= len(push) + 1  # N itemized figures + 1 total
        *itemized, printed_total = tokens
        assert len(itemized) == len(push)
        assert round(sum(itemized), 1) == pytest.approx(round(printed_total, 1), abs=0.15)

    def test_every_push_channel_named_in_slide_text(self):
        data = _logistics_plan()
        deck_kb = ppt._load_deck_kb()
        prs = _new_prs()
        ppt._build_slide_push_meets_pull(prs, data, deck_kb)
        assert len(prs.slides) == 1
        blob = "\n".join(_all_slide_text(prs))
        for needle in ("Programmatic", "Regional Job Boards", "Global Job Boards",
                       "Niche", "Social Media"):
            assert needle in blob, f"push channel not itemized on slide: {needle!r}"

    def test_atria_push_card_also_foots(self):
        data = _healthcare_plan()
        push, pull = ppt._push_pull_channel_split(data)
        total_budget = sum(d for _, d in push) + sum(d for _, d in pull)
        line = ppt._push_pull_split_line(push, total_budget)
        tokens = _money_tokens(line)
        *itemized, printed_total = tokens
        assert len(itemized) == len(push)
        assert round(sum(itemized), 1) == pytest.approx(round(printed_total, 1), abs=0.15)


# ---------------------------------------------------------------------------
# 2. No fabricated deck-only metrics
# ---------------------------------------------------------------------------
class TestNoFabricatedMetrics:
    def test_comparison_slide_has_no_reach_multiplier_or_diversity_score(self):
        for data in (_logistics_plan(), _healthcare_plan()):
            prs = _new_prs()
            ppt._build_slide_comparison_timeline(prs, data)
            blob = "\n".join(_all_slide_text(prs))
            assert "Reach Multiplier" not in blob
            assert "Channel Diversity Score" not in blob

    def test_healthcare_complications_no_longer_claim_unsourced_churn(self):
        complications = ppt._get_complications("healthcare_medical")
        blob = "\n".join(complications)
        assert "18% higher churn vs. 2023" not in blob
        assert "18%" not in blob or "churn" not in blob.lower()


# ---------------------------------------------------------------------------
# 2b. 4-critical cluster: comparison slide must read the FINAL allocation,
# not the static INDUSTRY_ALLOC_PROFILES split
# (visual:atria#1 / strategy:atria#1 / consistency:atria#1 / copy:both#1)
# ---------------------------------------------------------------------------
class TestComparisonSlideReadsFinalAllocation:
    def _plan_with_reweighted_allocation(self) -> dict:
        # Real post-reweight programmatic share is 28.1% (84,336/300,000) --
        # materially different from the STATIC healthcare_medical
        # INDUSTRY_ALLOC_PROFILES programmatic_dsp share (22, re-normalized
        # to ~21% once only 6 of 8 profile channels are selected). Any slide
        # that reads ``channels[*]["pct"]`` without reconciling against
        # ``_budget_allocation`` will show the wrong ~21% instead of 28%.
        data = _healthcare_plan()
        data["_budget_allocation"]["channel_allocations"] = {
            "programmatic_dsp": {"dollar_amount": 84_336.0, "percentage": 28.1},
            "global_boards": {"dollar_amount": 79_400.0, "percentage": 26.5},
            "regional_boards": {"dollar_amount": 56_400.0, "percentage": 18.8},
            "niche_boards": {"dollar_amount": 35_000.0, "percentage": 11.7},
            "social_media": {"dollar_amount": 20_850.0, "percentage": 7.0},
            "employer_branding": {"dollar_amount": 24_000.0, "percentage": 8.0},
        }
        return data

    def test_programmatic_allocation_matches_budget_allocation_metadata(self):
        data = self._plan_with_reweighted_allocation()
        ba_channel_alloc = data["_budget_allocation"]["channel_allocations"]
        expected_pct = round(ba_channel_alloc["programmatic_dsp"]["percentage"])

        prs = _new_prs()
        ppt._build_slide_comparison_timeline(prs, data)
        blob = "\n".join(_all_slide_text(prs))

        assert f"{expected_pct}%" in blob
        # The static-profile value (~21%) must NOT be the one shown.
        assert "21%" not in blob

    def test_programmatic_allocation_row_ties_to_channel_strategy_slide(self):
        # The comparison slide's Programmatic Allocation must equal the same
        # figure slide 5 (Channel Strategy) shows for the identical plan.
        data = self._plan_with_reweighted_allocation()
        prs5 = _new_prs()
        ppt._build_slide_channel_strategy(prs5, data)
        blob5 = "\n".join(_all_slide_text(prs5))

        prs9 = _new_prs()
        ppt._build_slide_comparison_timeline(prs9, data)
        blob9 = "\n".join(_all_slide_text(prs9))

        assert "28%" in blob5
        assert "28%" in blob9


# ---------------------------------------------------------------------------
# 3. Role Breakdown salary sourcing
# ---------------------------------------------------------------------------
class TestRoleBreakdownSalary:
    def _gold_standard_with_salary(self) -> dict:
        return {
            "difficulty_framework": [
                {"role_title": "Memory Care Associate", "seniority_level": "entry", "channel_emphasis": "niche_boards", "complexity_score": 5, "budget_weight": 1.0},
                {"role_title": "Nurse", "seniority_level": "professional", "channel_emphasis": "programmatic_dsp", "complexity_score": 8.5, "budget_weight": 1.8},
                {"role_title": "Cook", "seniority_level": "entry", "channel_emphasis": "regional_boards", "complexity_score": 3.5, "budget_weight": 0.6},
                {"role_title": "Driver", "seniority_level": "entry", "channel_emphasis": "global_boards", "complexity_score": 4, "budget_weight": 0.8},
            ],
            "city_level_data": {
                "new york, ny": {
                    "per_role_salary": {
                        "Memory Care Associate": {"median": 46920, "confidence": "benchmark"},
                        "Nurse": {"median": 103500, "confidence": "benchmark"},
                        # Deliberately distinct from every other role's median
                        # here so this fixture exercises the plain "(est.)"
                        # path; the salary-collision "(est., shared band)"
                        # path has its own dedicated test below.
                        "Cook": {"median": 42000, "confidence": "estimated"},
                        "Driver": {"median": 55200, "confidence": "benchmark"},
                    }
                }
            },
        }

    def test_role_breakdown_median_salary_helper_reads_gold_standard(self):
        gold = self._gold_standard_with_salary()
        median, is_estimated = ppt._role_breakdown_median_salary(gold, "Nurse")
        assert median == 103500
        assert is_estimated is False

    def test_role_breakdown_median_salary_flags_estimated(self):
        gold = self._gold_standard_with_salary()
        median, is_estimated = ppt._role_breakdown_median_salary(gold, "Cook")
        assert median == 42000
        assert is_estimated is True

    def test_role_breakdown_median_salary_missing_role_returns_none(self):
        gold = self._gold_standard_with_salary()
        median, is_estimated = ppt._role_breakdown_median_salary(gold, "Sales")
        assert median is None
        assert is_estimated is False

    def test_role_breakdown_slide_renders_real_salaries_not_all_dashes(self):
        data = _healthcare_plan()
        data["_gold_standard"] = self._gold_standard_with_salary()
        # Deliberately do NOT set data["_enriched"] -- the old code path
        # sourced salaries from _enriched.salary_data, which is exactly the
        # "enrichment offline" condition that used to blank out every row.
        prs = _new_prs()
        ppt._build_slide_role_breakdown(prs, data)
        assert len(prs.slides) == 1
        blob = "\n".join(_all_slide_text(prs))
        assert "$103K" in blob or "$104K" in blob  # Nurse median (~$103,500)
        assert "(est.)" in blob  # Cook is the estimated-confidence row
        # Not every row is '--' anymore.
        assert blob.count("--") < 4
        # strategy:atria#5: Difficulty and Budget Weight columns are real
        # per-role data from the same difficulty_framework rows, not '--'.
        assert "8.5/10" in blob  # Nurse complexity_score
        assert "1.8x" in blob  # Nurse budget_weight

    def test_role_breakdown_shows_dash_when_gold_standard_has_no_salary_data(self):
        data = _healthcare_plan()
        data["_gold_standard"] = {
            "difficulty_framework": self._gold_standard_with_salary()["difficulty_framework"]
        }
        prs = _new_prs()
        ppt._build_slide_role_breakdown(prs, data)
        assert len(prs.slides) == 1
        blob = "\n".join(_all_slide_text(prs))
        assert "--" in blob

    def test_role_breakdown_flags_duplicate_estimated_salary(self):
        # visual:atria#2: an ESTIMATED (fallback) salary that collides
        # byte-for-byte with another role's displayed salary (e.g. an
        # entry-level role priced identically to a mid-tier one) must be
        # flagged as a shared/implausible band, not silently repeated.
        data = _healthcare_plan()
        gold = self._gold_standard_with_salary()
        gold["city_level_data"]["new york, ny"]["per_role_salary"]["Cook"] = {
            "median": 103500,
            "confidence": "estimated",
        }
        data["_gold_standard"] = gold
        prs = _new_prs()
        ppt._build_slide_role_breakdown(prs, data)
        blob = "\n".join(_all_slide_text(prs))
        assert "(est., shared band)" in blob


# ---------------------------------------------------------------------------
# 4. Next Steps interpolation
# ---------------------------------------------------------------------------
class TestNextStepsInterpolation:
    def test_interpolates_client_roles_locations_budget(self):
        data = _logistics_plan()
        base = ppt._load_deck_kb().get("next_steps") or []
        out = ppt._interpolate_next_steps(base, data)
        blob = "\n".join(out)
        assert "Manpower - Amerigas" in blob
        assert "CDL A Driver" in blob
        assert "$150,000" in blob
        assert "6 months" in blob

    def test_two_clients_render_different_next_steps(self):
        base = ppt._load_deck_kb().get("next_steps") or []
        logistics_out = "\n".join(ppt._interpolate_next_steps(base, _logistics_plan()))
        healthcare_out = "\n".join(ppt._interpolate_next_steps(base, _healthcare_plan()))
        assert logistics_out != healthcare_out
        assert "Manpower - Amerigas" not in healthcare_out
        assert "Atria Senior Living" not in logistics_out

    def test_case_study_next_steps_slide_contains_client_name(self):
        data = _logistics_plan()
        deck_kb = ppt._load_deck_kb()
        prs = _new_prs()
        ppt._build_slide_case_study_next_steps(prs, data, deck_kb)
        assert len(prs.slides) == 1
        blob = "\n".join(_all_slide_text(prs))
        assert "Manpower - Amerigas" in blob

    def test_empty_steps_returns_empty(self):
        assert ppt._interpolate_next_steps([], _logistics_plan()) == []


# ---------------------------------------------------------------------------
# 5. Competitor vertical classification + tag
# ---------------------------------------------------------------------------
class TestCompetitorVerticalClassification:
    def test_amazon_is_talent_market_for_senior_living(self):
        assert ppt._classify_competitor_vertical("Amazon", "healthcare_medical") == "talent_market"

    def test_brookdale_is_industry_for_senior_living(self):
        assert ppt._classify_competitor_vertical("Brookdale Senior Living", "healthcare_medical") == "industry"

    def test_fedex_ups_are_industry_for_logistics(self):
        assert ppt._classify_competitor_vertical("FedEx", "logistics_supply_chain") == "industry"
        assert ppt._classify_competitor_vertical("UPS", "logistics_supply_chain") == "industry"
        assert ppt._classify_competitor_vertical("XPO Logistics", "logistics_supply_chain") == "industry"

    def test_amazon_is_industry_for_retail(self):
        assert ppt._classify_competitor_vertical("Amazon", "retail_consumer") == "industry"

    def test_national_carriers_are_talent_market_for_senior_living(self):
        # fix/plan-quality-8 (prod Atria bundle): UPS/FedEx are logistics
        # carriers, not same-vertical senior-living operators -- they only
        # compete for the same wage band/labor pool as a healthcare/senior-
        # living client, not the same residents/customers.
        for name in ("UPS", "FedEx", "Amazon", "DHL", "USPS", "Walmart", "Target"):
            assert (
                ppt._classify_competitor_vertical(name, "healthcare_medical")
                == "talent_market"
            ), name

    def test_genuine_senior_living_operators_stay_industry(self):
        for name in (
            "Brookdale Senior Living",
            "Sunrise Senior Living",
            "Atria Senior Living",
            "Sonida Senior Living",
            "Five Star Senior Living",
        ):
            assert (
                ppt._classify_competitor_vertical(name, "healthcare_medical")
                == "industry"
            ), name

    def test_competitive_landscape_slide_tags_amazon_talent_market(self):
        data = _healthcare_plan()  # competitors: Brookdale, Sunrise, Amazon
        prs = _new_prs()
        ppt._build_slide_competitive_landscape(prs, data)
        assert len(prs.slides) == 1
        blob = "\n".join(_all_slide_text(prs))
        assert "Talent-market competitor" in blob
        assert "Industry competitor" in blob

    def test_why_line_varies_by_vertical_type(self):
        industry_why = ppt._compose_competitor_why(
            "Brookdale Senior Living", {}, {"role": "Nurse", "city": "New York, NY", "vertical_type": "industry"}, 0
        )
        talent_market_why = ppt._compose_competitor_why(
            "Amazon", {}, {"role": "Nurse", "city": "New York, NY", "vertical_type": "talent_market"}, 0
        )
        assert industry_why != talent_market_why
        assert "wage" in talent_market_why.lower() or "labor pool" in talent_market_why.lower()


# ---------------------------------------------------------------------------
# 6. Company/competitor description truncation (word boundary, single ellipsis)
# ---------------------------------------------------------------------------
class TestTruncWord:
    def test_long_text_cuts_at_word_boundary_not_mid_word(self):
        desc = (
            "Atria Senior Living, Inc., commonly also known as Atria Senior "
            "Living, is a subsidiary of Atria Communities, Inc. and is a "
            "private company based in Louisville, Kentucky that operates "
            "senior living communities across the United States and Canada."
        )
        out = ppt._trunc_word(desc, 100)
        assert out.endswith("...")
        stem = out[: -len("...")]
        # Every word in the truncated stem must be a full word taken
        # verbatim from the source string -- i.e. the cut point was a
        # space in the original text, never inside a word.
        assert desc.startswith(stem)
        assert stem == "" or desc[len(stem)] == " " or len(stem) == len(desc)

    def test_no_double_ellipsis_when_source_already_truncated(self):
        # Simulates a pre-truncated upstream API description (already ends
        # in its own raw "...") -- re-truncating must not double the marker.
        pre_truncated = (
            "Atria Senior Living, Inc., also known as Atria Senior Living, "
            "is a subsidiary of Atria Communities, Inc. and offers senior "
            "living services nationwide..."
        )
        out = ppt._trunc_word(pre_truncated, 60)
        assert "......" not in out
        assert out.count("...") == 1
        assert out.endswith("...")

    def test_unicode_ellipsis_normalized(self):
        out = ppt._trunc_word("Atria Senior Living is a great company…", 500)
        assert "…" not in out
        assert not out.endswith("...")  # nothing was actually truncated

    def test_short_text_is_unchanged(self):
        assert ppt._trunc_word("Short description.", 100) == "Short description."

    def test_no_trailing_comma_before_ellipsis(self):
        # A cut that lands right after a comma-terminated word must not
        # leave a dangling "word,..." -- the comma is stripped too.
        desc = "Alpha, Beta, Gamma, Delta, Epsilon, Zeta, Eta, Theta, Iota"
        out = ppt._trunc_word(desc, 12)
        assert not re.search(r"[,;:-]\.\.\.$", out)


# ---------------------------------------------------------------------------
# 7. Competitor card overflow (cap to what fits above the footer)
# ---------------------------------------------------------------------------
class TestCompetitorCardOverflowFix:
    def test_four_plus_competitors_do_not_overflow_past_footer(self):
        data = _healthcare_plan()
        data["competitors"] = [
            "Brookdale Senior Living",
            "Sunrise Senior Living",
            "Sonida Senior Living",
            "Five Star Senior Living",
        ]
        prs = _new_prs()
        ppt._build_slide_competitive_landscape(prs, data)
        assert len(prs.slides) == 1
        slide = prs.slides[0]

        footer_rule_y = ppt.Inches(7.12)
        right_left = ppt.Inches(6.5)
        right_w = ppt.Inches(6.3)
        card_h = ppt.Inches(1.3)
        tol = ppt.Inches(0.05)

        # Identify competitor CARD shapes specifically (the white rounded-
        # rect card background, right_w wide x 1.3in tall) -- not the
        # "COMPETITOR LANDSCAPE" section header textbox, which shares the
        # same left/width but a much shorter height.
        card_bottoms = []
        for shape in slide.shapes:
            if (
                shape.left is not None
                and abs(shape.left - right_left) < tol
                and shape.width is not None
                and abs(shape.width - right_w) < tol
                and shape.height is not None
                and abs(shape.height - card_h) < tol
            ):
                card_bottoms.append(shape.top + shape.height)

        assert card_bottoms, "expected at least one competitor card shape"
        assert len(card_bottoms) <= 3, "no more than 3 competitor cards should render"
        for bottom in card_bottoms:
            assert bottom <= footer_rule_y, (
                f"competitor card bottom {bottom} exceeds footer y {footer_rule_y}"
            )

    def test_card_count_capped_at_three(self):
        data = _healthcare_plan()
        data["competitors"] = ["Brookdale Senior Living", "Sunrise Senior Living", "Sonida Senior Living", "Five Star Senior Living"]
        prs = _new_prs()
        ppt._build_slide_competitive_landscape(prs, data)
        blob = "\n".join(_all_slide_text(prs))
        # The 4th competitor should never appear -- a clean cap, not a
        # partially-drawn card.
        assert "Five Star Senior Living" not in blob


if __name__ == "__main__":
    sys.exit(pytest.main([__file__, "-v"]))
