"""Tests for the S89 PPTX deck polish (ppt_generator).

Covers the three executive-grade polish items, all offline (no network / no LLM
/ no Supabase -- the budget engine runs purely on its deterministic math):

  1. Currency-correctness -- money renders in the plan's own symbol (USD default,
     GBP / EUR / INR for non-US plans) via plan_currency wiring.
  2. Brand chart fonts -- bundled Poppins .ttf register with matplotlib, with a
     graceful DejaVu fallback when the fonts are absent.
  3. Channel table has a totals row.

Plus: the deck still builds end-to-end from a minimal dict.

(Item "Data-sources slide freshness" was retired 2026-08-04: its target,
_build_slide_data_sources, was dead code -- never called from generate_pptx
since the S50 slide-order change -- and was deleted along with this test in
the dead-slide-builder audit. See TestDataSourcesSlide in git history.)

Runs under pytest, or standalone: ``python3 tests/test_ppt_polish.py``.
"""

from __future__ import annotations

import datetime
import io
import sys
from pathlib import Path
from unittest import mock

import pytest

PROJECT_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(PROJECT_ROOT))

import ppt_generator as ppt  # noqa: E402
from pptx import Presentation  # noqa: E402


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def _build_alloc(country: str = "United States") -> dict:
    """Drive the real budget engine to produce an authentic _budget_allocation."""
    import budget_engine

    return budget_engine.calculate_budget_allocation(
        total_budget=150_000,
        roles=[
            {"title": "Registered Nurse", "count": 40, "tier": "mid"},
            {"title": "ICU Nurse", "count": 15, "tier": "senior"},
        ],
        locations=[{"city": "Metro", "state": "", "country": country}],
        industry="healthcare",
        channel_percentages={"Indeed": 40, "LinkedIn": 35, "Google Search Ads": 25},
        collar_type="white",
        campaign_start_month=9,
    )


def _plan(country: str = "United States") -> dict:
    return {
        "client_name": "Mercy Health",
        "industry": "healthcare",
        "budget": "$150,000",
        "locations": [country],
        "roles": ["Registered Nurse", "ICU Nurse"],
        "target_roles": [
            {"title": "Registered Nurse"},
            {"title": "ICU Nurse"},
        ],
        "_budget_allocation": _build_alloc(country),
    }


def _all_text(pptx_bytes: bytes) -> list[str]:
    prs = Presentation(io.BytesIO(pptx_bytes))
    out: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                out.append(shape.text_frame.text)
    return out


# ---------------------------------------------------------------------------
# 1. Currency-correctness
# ---------------------------------------------------------------------------
class TestCurrencyResolution:
    def test_defaults_to_usd_when_unknown(self):
        assert ppt._plan_currency_code({}) == "USD"
        assert ppt._plan_currency_code(None) == "USD"
        assert ppt._plan_currency_code({"locations": []}) == "USD"

    def test_resolves_country_from_locations(self):
        assert ppt._plan_currency_code({"locations": ["United Kingdom"]}) == "GBP"
        assert ppt._plan_currency_code({"locations": ["Germany"]}) == "EUR"
        assert ppt._plan_currency_code({"locations": ["Mumbai, India"]}) == "INR"

    def test_resolves_dict_locations(self):
        data = {"locations": [{"city": "London", "country": "United Kingdom"}]}
        assert ppt._plan_currency_code(data) == "GBP"

    def test_explicit_currency_code_wins(self):
        assert ppt._plan_currency_code({"currency_code": "eur"}) == "EUR"
        assert (
            ppt._plan_currency_code({"currency": "JPY", "locations": ["United States"]})
            == "JPY"
        )

    def test_set_active_currency_stashes_on_data(self):
        data = {"locations": ["United Kingdom"]}
        code = ppt._set_active_currency(data)
        assert code == "GBP"
        assert data["_plan_currency_code"] == "GBP"
        # restore default for other tests
        ppt._set_active_currency({})


class TestUsOnlyCampaignDetection:
    """_is_us_only_campaign gates whether _selected_channels() strips
    apac_regional / emea_regional. It must not fall back to "assume domestic"
    for countries the hardcoded substring blocklist doesn't list -- it should
    defer to plan_currency's country table, which already covers this."""

    def test_new_zealand_via_locations_and_country_is_not_us_only(self):
        # Exact repro (BUNDLE_QC_FINDINGS_2026-07-03.json slide-6 finding): a
        # NZ plan was silently losing apac_regional/emea_regional because
        # "new zealand" / "nz" / "auckland" weren't in the blocklist and
        # data["country"] was never even consulted.
        data = {"locations": ["Auckland, New Zealand"], "country": "New Zealand"}
        assert ppt._is_us_only_campaign(data) is False

    def test_new_zealand_via_country_field_only(self):
        assert ppt._is_us_only_campaign({"country": "New Zealand"}) is False

    def test_new_zealand_via_locations_only(self):
        assert (
            ppt._is_us_only_campaign({"locations": ["Auckland, New Zealand"]}) is False
        )

    def test_countries_missing_from_hardcoded_blocklist_are_detected(self):
        # None of these are in the substring blocklist, but plan_currency's
        # country table resolves all of them to a non-USD currency.
        for loc in (
            "Manila, Philippines",
            "Warsaw, Poland",
            "Dublin, Ireland",
            "Johannesburg, South Africa",
            "Seoul, South Korea",
            "Jakarta, Indonesia",
            "Lisbon, Portugal",
            "Zurich, Switzerland",
            "Moscow, Russia",
            "Stockholm, Sweden",
        ):
            assert ppt._is_us_only_campaign({"locations": [loc]}) is False, loc

    def test_plain_us_plan_is_still_us_only(self):
        assert ppt._is_us_only_campaign({"locations": ["United States"]}) is True
        assert (
            ppt._is_us_only_campaign(
                {"locations": ["San Francisco, CA", "New York, NY"]}
            )
            is True
        )
        assert ppt._is_us_only_campaign({"locations": ["Dallas, TX"]}) is True
        assert ppt._is_us_only_campaign({"locations": ["Dallas"]}) is True

    def test_indianapolis_does_not_false_positive_on_india_substring(self):
        assert ppt._is_us_only_campaign({"locations": ["Indianapolis, IN"]}) is True

    def test_no_locations_or_country_defaults_domestic(self):
        assert ppt._is_us_only_campaign({}) is True
        assert ppt._is_us_only_campaign({"locations": []}) is True

    def test_dict_locations_resolve_via_country_key(self):
        assert (
            ppt._is_us_only_campaign(
                {"locations": [{"city": "London", "country": "United Kingdom"}]}
            )
            is False
        )
        # Dict entries without a country/location key don't falsely trip
        # international -- matches _plan_currency_code's existing behavior.
        assert (
            ppt._is_us_only_campaign(
                {"locations": [{"city": "Austin", "state": "TX"}, "Remote"]}
            )
            is True
        )

    def test_target_region_short_circuit_still_wins(self):
        assert ppt._is_us_only_campaign({"target_region": "us_only"}) is True
        assert ppt._is_us_only_campaign({"target_region": "apac"}) is False
        # Explicit us_only overrides what locations would otherwise resolve to
        assert (
            ppt._is_us_only_campaign(
                {"target_region": "us_only", "locations": ["London, UK"]}
            )
            is True
        )

    def test_selected_channels_keeps_intl_channels_for_nz_plan(self):
        # End-to-end: the NZ repro must not strip apac_regional/emea_regional.
        data = {
            "locations": ["Auckland, New Zealand"],
            "country": "New Zealand",
            "channel_categories": {"apac_regional": True, "emea_regional": True},
        }
        selected = ppt._selected_channels(data)
        assert "apac_regional" in selected
        assert "emea_regional" in selected

    def test_selected_channels_strips_intl_channels_for_us_plan(self):
        data = {
            "locations": ["Dallas, TX"],
            "channel_categories": {"apac_regional": True, "emea_regional": True},
        }
        selected = ppt._selected_channels(data)
        assert "apac_regional" not in selected
        assert "emea_regional" not in selected


class TestCurrencyFormatting:
    def teardown_method(self):
        ppt._set_active_currency({})  # reset to USD between tests

    def test_fmt_currency_uses_active_currency(self):
        ppt._set_active_currency({"locations": ["United Kingdom"]})
        assert ppt._fmt_currency(150000) == "£150,000"
        assert ppt._fmt_currency(1_200_000, compact=True) == "£1.2M"

    def test_fmt_currency_explicit_override_is_usd(self):
        ppt._set_active_currency({"locations": ["United Kingdom"]})
        # benchmark figures opt back into USD explicitly
        assert ppt._fmt_currency(35, currency="USD") == "$35"

    def test_format_budget_display_localizes(self):
        ppt._set_active_currency({"locations": ["Germany"]})
        assert ppt._format_budget_display("$150,000") == "€150K"

    def test_format_salary_localizes(self):
        ppt._set_active_currency({"locations": ["India"]})
        assert ppt._format_salary(85000) == "₹85K"

    def test_none_and_default_unchanged(self):
        ppt._set_active_currency({})
        assert ppt._fmt_currency(None) == "—"
        assert ppt._fmt_currency(150000) == "$150,000"


class TestNonUsdDeckRendering:
    def test_uk_plan_renders_pounds_not_dollars(self):
        pptx_bytes = ppt.generate_pptx(_plan("United Kingdom"))
        texts = _all_text(pptx_bytes)
        # The plan's own money figures use the pound symbol...
        pound_cells = [t for t in texts if t.startswith("£")]
        assert pound_cells, "expected GBP-formatted money in a UK plan"
        # ...and the localized total investment hero is present.
        assert any("£150" in t for t in texts)
        ppt._set_active_currency({})  # reset

    def test_us_plan_still_uses_dollars(self):
        pptx_bytes = ppt.generate_pptx(_plan("United States"))
        texts = _all_text(pptx_bytes)
        assert any(t.startswith("$") and "," in t for t in texts)


# ---------------------------------------------------------------------------
# 2. Brand chart fonts
# ---------------------------------------------------------------------------
class TestChartFonts:
    def test_bundled_poppins_files_exist(self):
        fonts_dir = PROJECT_ROOT / "fonts"
        for name in (
            "Poppins-Regular.ttf",
            "Poppins-SemiBold.ttf",
            "Poppins-Bold.ttf",
        ):
            assert (fonts_dir / name).is_file(), f"missing bundled font {name}"

    def test_register_returns_poppins_when_present(self):
        if not ppt._HAS_MATPLOTLIB:
            pytest.skip("matplotlib not installed")
        assert ppt._register_chart_fonts() == "Poppins"

    def test_register_falls_back_to_dejavu_when_dir_missing(self):
        # Point the resolver at a directory with no fonts -> graceful fallback.
        with mock.patch.object(ppt, "_FONTS_DIR", PROJECT_ROOT / "no_such_fonts_dir"):
            assert ppt._register_chart_fonts() == "DejaVu Sans"

    def test_pie_chart_still_renders(self):
        if not ppt._HAS_MATPLOTLIB:
            pytest.skip("matplotlib not installed")
        png = ppt._generate_pie_chart_image(["Indeed", "LinkedIn"], [60.0, 40.0])
        assert png and png[:4] == b"\x89PNG"


# ---------------------------------------------------------------------------
# 3. Channel table totals row
# ---------------------------------------------------------------------------
class TestTotalsRow:
    def teardown_method(self):
        ppt._set_active_currency({})

    def test_budget_slide_has_totals_row(self):
        data = _plan("United States")
        ppt._set_active_currency(data)
        prs = Presentation()
        prs.slide_width = ppt.SLIDE_WIDTH
        prs.slide_height = ppt.SLIDE_HEIGHT
        ppt._build_slide_budget_allocation(prs, data)
        texts = []
        for shape in prs.slides[0].shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        assert any(t.strip() == "Total" for t in texts), "missing totals row"


# ---------------------------------------------------------------------------
# End-to-end: deck still builds from a minimal dict
# ---------------------------------------------------------------------------
class TestEndToEnd:
    def teardown_method(self):
        ppt._set_active_currency({})

    def test_minimal_dict_builds_valid_pptx(self):
        data = {
            "client_name": "Acme Co",
            "industry": "healthcare",
            "budget": "$100,000",
            "locations": ["New York, NY"],
            "roles": ["Registered Nurse"],
        }
        pptx_bytes = ppt.generate_pptx(data)
        assert isinstance(pptx_bytes, bytes) and len(pptx_bytes) > 10_000
        # Valid .pptx is a zip archive (PK magic bytes).
        assert pptx_bytes[:2] == b"PK"
        prs = Presentation(io.BytesIO(pptx_bytes))
        assert len(prs.slides) >= 5

    def test_full_sample_shape_builds(self):
        pptx_bytes = ppt.generate_pptx(_plan("United States"))
        assert pptx_bytes[:2] == b"PK"


# ---------------------------------------------------------------------------
# 4. S6 -- channel percentages always reconcile to 100%, and the slide-6
#    breakdown table always foots (visible rows sum to the printed Total),
#    even when there are more channels than fit in the visible row cap.
# ---------------------------------------------------------------------------
class TestChannelPercentageFooting:
    def test_largest_remainder_sums_to_100_on_rounding_drift(self):
        # Independently rounding each of these would give 31+21+16+13+11+4+3+2=101
        # (the exact drift confirmed on the Pratt & Whitney NZ deck).
        values = {
            "a": 30.9,
            "b": 21.0,
            "c": 16.0,
            "d": 12.6,
            "e": 10.9,
            "f": 3.8,
            "g": 3.0,
            "h": 2.0,
        }
        result = ppt._largest_remainder_round(values)
        assert sum(result.values()) == 100

    def test_largest_remainder_sums_to_100_when_subset_is_incomplete(self):
        # If the caller only has a subset of the full channel mix (raw values
        # sum to 93, not ~100), reconciliation must still land on exactly 100
        # rather than under-shooting to 99 (each item capped at +1).
        values = {
            "a": 23.0,
            "b": 18.0,
            "c": 33.0,
            "d": 3.0,
            "e": 13.0,
            "f": 3.0,
        }
        result = ppt._largest_remainder_round(values)
        assert sum(result.values()) == 100

    def test_largest_remainder_empty_and_zero_safe(self):
        assert ppt._largest_remainder_round({}) == {}
        result = ppt._largest_remainder_round({"a": 0, "b": 0})
        assert result == {"a": 0, "b": 0}

    def _render_budget_slide_text(
        self, data: dict, channels_selected: dict
    ) -> list[str]:
        orig = ppt._selected_channels
        ppt._selected_channels = lambda _data: {
            k: dict(v) for k, v in channels_selected.items()
        }
        try:
            prs = Presentation()
            prs.slide_width = ppt.SLIDE_WIDTH
            prs.slide_height = ppt.SLIDE_HEIGHT
            ppt._build_slide_budget_allocation(prs, data)
        finally:
            ppt._selected_channels = orig
        return [
            shape.text_frame.text.strip()
            for shape in prs.slides[0].shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]

    def test_table_foots_with_rollup_row_when_more_than_six_channels(self):
        from pptx.dml.color import RGBColor

        clr = RGBColor(0x5A, 0x54, 0xBE)
        ba_channel_alloc = {
            "niche_boards": {
                "label": "Niche / Industry Boards",
                "percentage": 33.0,
                "dollar_amount": 49500,
                "projected_applications": 2600,
                "projected_hires": 11,
                "cpa": 19,
            },
            "programmatic_dsp": {
                "label": "Programmatic DSP",
                "percentage": 23.0,
                "dollar_amount": 34500,
                "projected_applications": 1600,
                "projected_hires": 7,
                "cpa": 22,
            },
            "global_boards": {
                "label": "Global Job Boards",
                "percentage": 18.0,
                "dollar_amount": 27000,
                "projected_applications": 1400,
                "projected_hires": 6,
                "cpa": 19,
            },
            "regional_boards": {
                "label": "Regional Boards",
                "percentage": 13.0,
                "dollar_amount": 19500,
                "projected_applications": 1000,
                "projected_hires": 4,
                "cpa": 19,
            },
            "apac_regional": {
                "label": "APAC Regional",
                "percentage": 5.8,
                "dollar_amount": 8700,
                "projected_applications": 100,
                "projected_hires": 1,
                "cpa": 87,
            },
            "social_media": {
                "label": "Social Media",
                "percentage": 3.0,
                "dollar_amount": 4500,
                "projected_applications": 22,
                "projected_hires": 0,
                "cpa": 205,
            },
            "employer_branding": {
                "label": "Employer Branding",
                "percentage": 3.0,
                "dollar_amount": 4500,
                "projected_applications": 28,
                "projected_hires": 0,
                "cpa": 161,
            },
            "emea_regional": {
                "label": "EMEA Regional",
                "percentage": 1.0,
                "dollar_amount": 1500,
                "projected_applications": 15,
                "projected_hires": 0,
                "cpa": 100,
            },
        }
        channels_selected = {
            k: {"label": v["label"], "pct": round(v["percentage"]), "color": clr}
            for k, v in ba_channel_alloc.items()
        }
        data = {
            "client_name": "Pratt & Whitney New Zealand",
            "industry": "aerospace_defense",
            "budget": "$150,000",
            "_budget_allocation": {
                "channel_allocations": ba_channel_alloc,
                "total_projected": {
                    "applications": 6765,
                    "hires": 29,
                    "cost_per_application": 22,
                    "cost_per_hire": 5200,
                },
                "metadata": {"total_budget": 150000},
            },
        }
        texts = self._render_budget_slide_text(data, channels_selected)

        # Rollup row present for the 3 channels that don't fit in 6 visible rows.
        assert any("smaller channels" in t for t in texts), "missing rollup row"

        # Extract percentage cells (strings ending in "%") in table order:
        # header/KPI cards also contain "%"-free text, so filter on the pattern.
        pct_cells = [
            t for t in texts if t.endswith("%") and t[:-1].replace(".", "").isdigit()
        ]
        # Last one is the Total row's percentage; the rest are visible-row pcts.
        assert pct_cells, "no percentage cells found"
        total_pct = int(pct_cells[-1].rstrip("%"))
        visible_pct_sum = sum(int(p.rstrip("%")) for p in pct_cells[:-1])
        assert total_pct == 100, f"Total row must print 100%, got {total_pct}%"
        assert visible_pct_sum == total_pct, (
            f"visible rows ({visible_pct_sum}%) must foot to the printed "
            f"Total ({total_pct}%)"
        )

    def test_table_no_rollup_when_six_or_fewer_channels(self):
        from pptx.dml.color import RGBColor

        clr = RGBColor(0x5A, 0x54, 0xBE)
        ba_channel_alloc = {
            "niche_boards": {
                "label": "Niche / Industry Boards",
                "percentage": 40.0,
                "dollar_amount": 60000,
                "projected_applications": 3000,
                "projected_hires": 13,
                "cpa": 20,
            },
            "programmatic_dsp": {
                "label": "Programmatic DSP",
                "percentage": 30.0,
                "dollar_amount": 45000,
                "projected_applications": 2000,
                "projected_hires": 9,
                "cpa": 22,
            },
            "global_boards": {
                "label": "Global Job Boards",
                "percentage": 30.0,
                "dollar_amount": 45000,
                "projected_applications": 2000,
                "projected_hires": 8,
                "cpa": 22,
            },
        }
        channels_selected = {
            k: {"label": v["label"], "pct": round(v["percentage"]), "color": clr}
            for k, v in ba_channel_alloc.items()
        }
        data = {
            "client_name": "Acme Co",
            "industry": "healthcare",
            "budget": "$150,000",
            "_budget_allocation": {
                "channel_allocations": ba_channel_alloc,
                "total_projected": {
                    "applications": 7000,
                    "hires": 30,
                    "cost_per_application": 21,
                    "cost_per_hire": 5000,
                },
                "metadata": {"total_budget": 150000},
            },
        }
        texts = self._render_budget_slide_text(data, channels_selected)
        assert not any(
            "smaller channels" in t for t in texts
        ), "should not add a rollup row when all channels already fit"


if __name__ == "__main__":
    sys.exit(pytest.main([__file__, "-q"]))
