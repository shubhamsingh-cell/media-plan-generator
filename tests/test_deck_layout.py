"""Layout-integrity regression tests for the generated PPTX deck.

These lock in the June 2026 "baseline polish" pass that fixed the deck's
overlay/spacing/overflow defects, so they can't silently regress:

  1. Text-bearing shapes stay within the slide bounds (no off-canvas KPIs /
     benchmark-table rows). Decorative fills (ovals/rectangles with no text) may
     intentionally bleed off the edge.
  2. No text renders below an 8pt readability floor.
  3. The brand font (Poppins) is embedded so the deck is on-brand on machines
     that don't have it installed.

Offline: drives the real budget engine (deterministic math, no network/LLM).
Runs under pytest, or standalone: ``python3 tests/test_deck_layout.py``.
"""

from __future__ import annotations

import io
import sys
import zipfile
from pathlib import Path

PROJECT_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(PROJECT_ROOT))

import ppt_generator as ppt  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.util import Emu  # noqa: E402

EMU_PER_IN = 914400
TOL = Emu(int(0.06 * EMU_PER_IN))  # 0.06in tolerance for rounding


def _worst_case_plan() -> dict:
    """A content-heavy plan that exercises the layouts that used to overflow:
    6 channels (KPI strip + attribution band), a full budget allocation, a
    creative-quality badge, and many benchmark rows."""
    import budget_engine

    alloc = budget_engine.calculate_budget_allocation(
        total_budget=375_000,
        roles=[
            {"title": "Software Engineer", "count": 40, "tier": "senior"},
            {"title": "Data Scientist", "count": 20, "tier": "senior"},
            {"title": "DevOps Engineer", "count": 10, "tier": "mid"},
            {"title": "Product Manager", "count": 5, "tier": "senior"},
        ],
        locations=[{"city": "New York", "state": "NY", "country": "United States"}],
        industry="technology_engineering",
        channel_percentages={
            "Programmatic DSP": 36,
            "Niche / Industry Boards": 21,
            "Social Media": 18,
            "Global Job Boards": 16,
            "Regional Boards": 6,
            "Employer Branding": 3,
        },
        collar_type="white",
        campaign_start_month=6,
    )
    return {
        "client_name": "Uber",
        "industry": "technology_engineering",
        "industry_label": "Technology & Engineering",
        "locations": ["New York, NY"],
        "roles": [
            "Software Engineer",
            "Data Scientist",
            "DevOps Engineer",
            "Product Manager",
        ],
        "budget": "$375,000",
        "work_environment": "hybrid",
        "channel_categories": {
            "regional_boards": True,
            "global_boards": True,
            "niche_boards": True,
            "social_media": True,
            "programmatic_dsp": True,
            "employer_branding": True,
        },
        "_budget_allocation": alloc,
        "_synthesized": {
            "job_market_demand": {
                "Software Engineer": {
                    "total_postings": 150000,
                    "avg_salary": 170000,
                    "market_temperature": "hot",
                },
                "Data Scientist": {
                    "total_postings": 45000,
                    "avg_salary": 165000,
                    "market_temperature": "hot",
                },
            },
            "ad_platform_analysis": {
                "google_ads": {
                    "platform_name": "Google Ads",
                    "CPC": 4.2,
                    "CPA": 48,
                    "fit_score": 90,
                },
                "meta": {
                    "platform_name": "Meta",
                    "CPC": 2.1,
                    "CPA": 33,
                    "fit_score": 50,
                },
                "linkedin": {
                    "platform_name": "LinkedIn Ads",
                    "CPC": 9.2,
                    "CPA": 120,
                    "fit_score": 100,
                },
                "tiktok": {
                    "platform_name": "TikTok Ads",
                    "CPC": 1.5,
                    "CPA": 22,
                    "fit_score": 20,
                },
                "bing": {
                    "platform_name": "Bing Ads",
                    "CPC": 3.5,
                    "CPA": 56,
                    "fit_score": 60,
                },
            },
        },
        "_creative_quality_score": {"score": 45, "grade": "F"},
    }


def _deck() -> bytes:
    return ppt.generate_pptx(_worst_case_plan())


def _iter_shapes(shapes):
    for sh in shapes:
        yield sh
        if sh.shape_type == 6:  # group
            try:
                yield from _iter_shapes(sh.shapes)
            except Exception:
                pass


def test_text_shapes_within_slide_bounds():
    """No text-bearing shape may extend past the slide edge. Decorative fills
    (no text) are allowed to bleed (full-bleed cover accents)."""
    prs = Presentation(io.BytesIO(_deck()))
    sw, sh = prs.slide_width, prs.slide_height
    offenders = []
    for idx, slide in enumerate(prs.slides, 1):
        for shape in _iter_shapes(slide.shapes):
            if not shape.has_text_frame or not shape.text_frame.text.strip():
                continue
            try:
                l, t, w, h = shape.left, shape.top, shape.width, shape.height
            except Exception:
                continue
            if None in (l, t, w, h):
                continue
            if l < -TOL or t < -TOL or (l + w) > sw + TOL or (t + h) > sh + TOL:
                offenders.append(
                    f"slide {idx}: {shape.name!r} ({shape.text_frame.text[:30]!r}) "
                    f"L={l/EMU_PER_IN:.2f} R={(l+w)/EMU_PER_IN:.2f} "
                    f"B={(t+h)/EMU_PER_IN:.2f} (slide {sw/EMU_PER_IN:.2f}x{sh/EMU_PER_IN:.2f})"
                )
    assert not offenders, "Text off-slide:\n" + "\n".join(offenders)


def test_no_text_below_8pt():
    """Every text run must be >= 8pt (readability floor)."""
    prs = Presentation(io.BytesIO(_deck()))
    tiny = []
    for idx, slide in enumerate(prs.slides, 1):
        for shape in _iter_shapes(slide.shapes):
            if not shape.has_text_frame:
                continue
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    if (
                        run.font.size is not None
                        and run.font.size.pt < 8
                        and run.text.strip()
                    ):
                        tiny.append(
                            f"slide {idx}: {run.font.size.pt}pt {run.text[:30]!r}"
                        )
    assert not tiny, "Sub-8pt text:\n" + "\n".join(tiny)


def test_brand_font_embedded():
    """The deck must embed Poppins so it renders on-brand everywhere."""
    data = _deck()
    with zipfile.ZipFile(io.BytesIO(data)) as zf:
        names = zf.namelist()
        assert any(n.endswith(".fntdata") for n in names), "no embedded font parts"
        pres = zf.read("ppt/presentation.xml").decode("utf-8")
        assert 'embedTrueTypeFonts="1"' in pres, "embedTrueTypeFonts flag not set"


if __name__ == "__main__":
    test_text_shapes_within_slide_bounds()
    test_no_text_below_8pt()
    test_brand_font_embedded()
    print("All deck-layout regression tests passed.")
