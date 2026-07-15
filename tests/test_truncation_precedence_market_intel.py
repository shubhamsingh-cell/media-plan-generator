"""S95 follow-up (wave 2): `X or ""[:N]` / `X or [][:N]` operator-precedence
no-ops across the market/research modules, verified by targeted repro before
fixing. Companion to tests/test_truncation_precedence_fixes.py (which covered
excel_v2.py + data_synthesizer.py); this file covers the remaining 4 modules.

Python precedence: slicing binds tighter than `or`, so
``brief.get("competitors") or [][:5]`` slices the empty-list LITERAL and the
intended cap is silently a no-op -- the unsliced (or empty-fallback) value
flows through unbounded. Audited sites fixed here:

  1. competitive_intel.generate_competitive_brief -- competitor list fed to
     the LLM narrative prompt (`brief.get("competitors") or [][:5]`, comment
     "Cap for token budget"). The token-budget cap was a no-op; the brief's
     own `competitors` field (used elsewhere, e.g. comparison tables) stays
     unbounded by design -- only the narrative-prompt COPY is capped.
  2. competitive_intel.generate_competitive_excel footer
     (`brief.get('generated_at') or ''[:10]`) -- the ISO timestamp's 10-char
     date-prefix cap was a no-op; the full "YYYY-MM-DDTHH:MM:SS.ffffffZ"
     string rendered in the sheet subtitle.
  3. competitive_intel.generate_competitive_ppt footer -- same pattern as #2,
     on the title-slide subtitle text box.
  4. market_pulse.generate_pulse_email_html "Top 3 industries" section
     (`ind_data.get("industries") or [][:3]`) -- the digest-email industry
     table rendered all available industries instead of the top 3.
  5. market_intel_reports._build_executive_summary
     (`", ".join(report.get("locations") or [][:3]) or "multiple regions"`)
     -- the 3-location cap on the executive-summary sentence was a no-op;
     the "multiple regions" fallback for missing/empty locations already
     worked (an empty list joins to "", which is falsy either way).
  6. market_intel_reports.generate_intel_excel "Locations" info row
     (`", ".join(report_data.get("locations") or [][:5])`) -- 5-location cap
     was a no-op.
  7. market_intel_reports.generate_intel_ppt title-slide subtitle
     (`", ".join(report_data.get("locations") or [][:3])`) -- 3-location cap
     was a no-op.
  8. research.get_location_boards international branch
     (`for board_entry in country_boards.get("boards") or [][:5]`) -- the
     5-board cap was a no-op; repro with Germany (12 configured boards)
     showed all 12 flowing through unbounded before the fix.

All 8 sites are plain list/string slices (not prose mid-word truncation), so
parenthesizing the slice -- `(X or [])[:N]` / `(X or "")[:N]` -- is the
complete fix; no word-boundary helper is needed.

A source-level guard test keeps the precedence pattern from reappearing in
the four audited files.

Runs under pytest, or standalone:
``python3 tests/test_truncation_precedence_market_intel.py``.
"""

from __future__ import annotations

import io
import re
import sys
from pathlib import Path

PROJECT_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(PROJECT_ROOT))

import openpyxl  # noqa: E402
from pptx import Presentation  # noqa: E402

import competitive_intel  # noqa: E402
import market_intel_reports  # noqa: E402
import market_pulse  # noqa: E402
import research  # noqa: E402

_AUDITED_FILES = (
    "competitive_intel.py",
    "market_pulse.py",
    "market_intel_reports.py",
    "research.py",
)


# ===========================================================================
# Guard: the precedence pattern must not reappear in the audited files
# ===========================================================================
def test_slice_precedence_pattern_absent_from_audited_files():
    """`or ""[...]` / `or ''[...]` / `or [][...]` always slice the empty
    LITERAL -- any occurrence is a bug by construction."""
    pattern = re.compile(r"""or\s+(""|''|\[\])\s*\[""")
    for fname in _AUDITED_FILES:
        src = (PROJECT_ROOT / fname).read_text(encoding="utf-8")
        hits = [
            f"{fname}:{i}"
            for i, line in enumerate(src.splitlines(), 1)
            if pattern.search(line)
        ]
        assert not hits, f"slice-binds-to-literal precedence bug at: {hits}"


# ===========================================================================
# Shared helpers
# ===========================================================================
def _pptx_texts(pptx_bytes: bytes):
    prs = Presentation(io.BytesIO(pptx_bytes))
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                yield shape.text_frame.text


def _excel_strings(xlsx_bytes: bytes):
    wb = openpyxl.load_workbook(io.BytesIO(xlsx_bytes))
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            for val in row:
                if isinstance(val, str):
                    yield val


_MINIMAL_BRIEF = {
    "company": {"name": "Acme Corp", "industry": "Logistics"},
    "competitors": [],
    "comparison_matrix": [],
    "hiring_activity": {},
    "ad_benchmarks": {},
    "market_trends": {},
    "recommendations": [],
}


def _brief(**overrides) -> dict:
    b = dict(_MINIMAL_BRIEF)
    b.update(overrides)
    return b


# ===========================================================================
# Site 1: competitive_intel.generate_competitive_brief narrative-prompt cap
# ===========================================================================
def test_competitor_narrative_input_capped_at_five():
    """The LLM-prompt copy of `competitors` must be capped at 5 even though
    the brief's own `competitors` field (used for comparison tables etc.)
    stays unbounded -- only the token-budget copy is capped."""
    captured = {}
    original = competitive_intel._generate_competitive_narrative

    def _fake_narrative(competitor_data):
        captured["competitors"] = competitor_data.get("competitors")
        return "stub narrative"

    competitive_intel._generate_competitive_narrative = _fake_narrative
    try:
        analysis_results = {
            "competitor_analysis": {
                "company": {"name": "Acme"},
                "competitors": [{"name": f"Competitor {i}"} for i in range(8)],
                "comparison_matrix": [],
            },
            "hiring_activity": {},
            "ad_benchmarks": {},
            "market_trends": {},
        }
        brief = competitive_intel.generate_competitive_brief(analysis_results)
    finally:
        competitive_intel._generate_competitive_narrative = original

    assert len(captured["competitors"]) == 5, "the token-budget cap is still a no-op"
    assert captured["competitors"] == [{"name": f"Competitor {i}"} for i in range(5)]
    # The brief's own field is a separate, unbounded copy -- not affected by
    # the narrative-prompt cap.
    assert len(brief["competitors"]) == 8


def test_competitor_narrative_input_short_list_unchanged():
    captured = {}
    original = competitive_intel._generate_competitive_narrative

    def _fake_narrative(competitor_data):
        captured["competitors"] = competitor_data.get("competitors")
        return "stub narrative"

    competitive_intel._generate_competitive_narrative = _fake_narrative
    try:
        analysis_results = {
            "competitor_analysis": {
                "company": {"name": "Acme"},
                "competitors": [{"name": "OnlyCompetitor"}],
                "comparison_matrix": [],
            },
            "hiring_activity": {},
            "ad_benchmarks": {},
            "market_trends": {},
        }
        competitive_intel.generate_competitive_brief(analysis_results)
    finally:
        competitive_intel._generate_competitive_narrative = original

    assert captured["competitors"] == [{"name": "OnlyCompetitor"}]


def test_competitor_narrative_input_missing_falls_back_to_empty_list():
    captured = {}
    original = competitive_intel._generate_competitive_narrative

    def _fake_narrative(competitor_data):
        captured["competitors"] = competitor_data.get("competitors")
        return "stub narrative"

    competitive_intel._generate_competitive_narrative = _fake_narrative
    try:
        # No "competitor_analysis" key at all.
        competitive_intel.generate_competitive_brief(
            {"hiring_activity": {}, "ad_benchmarks": {}, "market_trends": {}}
        )
    finally:
        competitive_intel._generate_competitive_narrative = original

    assert captured["competitors"] == []


# ===========================================================================
# Sites 2 & 3: competitive_intel Excel/PPT footer date-prefix cap
# ===========================================================================
_LONG_TIMESTAMP = "2026-07-15T10:07:48.123456Z"


def test_excel_footer_timestamp_capped_to_date_prefix():
    xlsx_bytes = competitive_intel.generate_competitive_excel(
        _brief(generated_at=_LONG_TIMESTAMP), "Acme Corp"
    )
    hits = [v for v in _excel_strings(xlsx_bytes) if v.startswith("Generated ")]
    assert hits, "footer subtitle did not render"
    assert hits[0] == "Generated 2026-07-15 | Powered by Nova AI Suite", hits[0]
    assert _LONG_TIMESTAMP not in hits[0], "the 10-char cap is still a no-op"


def test_excel_footer_timestamp_short_value_unchanged():
    xlsx_bytes = competitive_intel.generate_competitive_excel(
        _brief(generated_at="2026-07-15"), "Acme Corp"
    )
    hits = [v for v in _excel_strings(xlsx_bytes) if v.startswith("Generated ")]
    assert hits[0] == "Generated 2026-07-15 | Powered by Nova AI Suite"


def test_excel_footer_timestamp_missing_falls_back_to_empty_prefix():
    xlsx_bytes = competitive_intel.generate_competitive_excel(
        _brief(generated_at=None), "Acme Corp"
    )
    hits = [v for v in _excel_strings(xlsx_bytes) if v.startswith("Generated ")]
    assert hits[0] == "Generated  | Powered by Nova AI Suite"


def test_ppt_footer_timestamp_capped_to_date_prefix():
    pptx_bytes = competitive_intel.generate_competitive_ppt(
        _brief(generated_at=_LONG_TIMESTAMP), "Acme Corp"
    )
    hits = [t for t in _pptx_texts(pptx_bytes) if t.startswith("Generated ")]
    assert hits, "footer text box did not render"
    assert hits[0] == "Generated 2026-07-15 | Powered by Nova AI Suite", hits[0]
    assert _LONG_TIMESTAMP not in hits[0], "the 10-char cap is still a no-op"


def test_ppt_footer_timestamp_missing_falls_back_to_empty_prefix():
    pptx_bytes = competitive_intel.generate_competitive_ppt(
        _brief(generated_at=None), "Acme Corp"
    )
    hits = [t for t in _pptx_texts(pptx_bytes) if t.startswith("Generated ")]
    assert hits[0] == "Generated  | Powered by Nova AI Suite"


# ===========================================================================
# Site 4: market_pulse.generate_pulse_email_html "Top 3 industries"
# ===========================================================================
def _pulse_report(industries):
    return {
        "report_date_display": "July 15, 2026",
        "period": "Q3 2026",
        "report_id": "abc123",
        "key_takeaways": [],
        "cpc_trends": {"available": False},
        "industry_spotlight": {"available": True, "industries": industries},
    }


def test_email_industry_table_capped_at_three():
    industries = [
        {"label": f"Industry {i}", "avg_cpc": 5.0, "avg_change_pct": 1.0}
        for i in range(6)
    ]
    html = market_pulse.generate_pulse_email_html(_pulse_report(industries))
    rendered = [f"Industry {i}" in html for i in range(6)]
    assert rendered[:3] == [True, True, True], "top-3 cap is still a no-op"
    assert rendered[3:] == [False, False, False], "extra industries leaked through"


def test_email_industry_table_short_list_unchanged():
    industries = [
        {"label": "OnlyIndustry", "avg_cpc": 5.0, "avg_change_pct": 1.0}
    ]
    html = market_pulse.generate_pulse_email_html(_pulse_report(industries))
    assert "OnlyIndustry" in html


def test_email_industry_table_empty_list_renders_no_rows():
    html = market_pulse.generate_pulse_email_html(_pulse_report([]))
    # No industry label text should appear; must not raise.
    assert "Industry 0" not in html


# ===========================================================================
# Site 5: market_intel_reports._build_executive_summary location cap
# ===========================================================================
def _mir_report(**overrides):
    report = {
        "industry": "logistics_supply_chain",
        "role_category": "cdl_driver",
        "locations": [],
        "market_overview": {"difficulty_score": 55, "growth_rate": 0.05},
    }
    report.update(overrides)
    return report


def test_executive_summary_locations_capped_at_three():
    report = _mir_report(locations=[f"City{i}, ST" for i in range(6)])
    summary = market_intel_reports._build_executive_summary(report)
    rendered = [f"City{i}" in summary for i in range(6)]
    assert rendered[:3] == [True, True, True], "3-location cap is still a no-op"
    assert rendered[3:] == [False, False, False], "extra locations leaked through"


def test_executive_summary_locations_short_list_unchanged():
    report = _mir_report(locations=["Dallas, TX", "Austin, TX"])
    summary = market_intel_reports._build_executive_summary(report)
    assert "Dallas, TX" in summary and "Austin, TX" in summary


def test_executive_summary_locations_empty_falls_back_to_multiple_regions():
    report = _mir_report(locations=[])
    summary = market_intel_reports._build_executive_summary(report)
    assert "multiple regions" in summary


def test_executive_summary_locations_missing_key_falls_back_to_multiple_regions():
    report = _mir_report()
    del report["locations"]
    summary = market_intel_reports._build_executive_summary(report)
    assert "multiple regions" in summary


# ===========================================================================
# Site 6: market_intel_reports.generate_intel_excel "Locations" cell
# ===========================================================================
def _intel_report(**overrides):
    report = {
        "industry": "logistics_supply_chain",
        "role_category": "cdl_driver",
        "locations": [],
        "report_metadata": {},
        "market_overview": {},
    }
    report.update(overrides)
    return report


def test_excel_locations_capped_at_five():
    report = _intel_report(locations=[f"City{i}, ST" for i in range(7)])
    xlsx_bytes = market_intel_reports.generate_intel_excel(report)
    hits = [v for v in _excel_strings(xlsx_bytes) if v.startswith("City0")]
    assert hits, "Locations cell did not render"
    locs_cell = hits[0]
    rendered = [f"City{i}" in locs_cell for i in range(7)]
    assert rendered[:5] == [True] * 5, "5-location cap is still a no-op"
    assert rendered[5:] == [False, False], "extra locations leaked through"


def test_excel_locations_short_list_unchanged():
    report = _intel_report(locations=["Dallas, TX", "Austin, TX"])
    xlsx_bytes = market_intel_reports.generate_intel_excel(report)
    hits = [v for v in _excel_strings(xlsx_bytes) if "Dallas, TX" in v]
    assert hits and "Austin, TX" in hits[0]


# ===========================================================================
# Site 7: market_intel_reports.generate_intel_ppt subtitle
# ===========================================================================
def test_ppt_locations_capped_at_three():
    report = _intel_report(locations=[f"City{i}, ST" for i in range(7)])
    pptx_bytes = market_intel_reports.generate_intel_ppt(report)
    hits = [t for t in _pptx_texts(pptx_bytes) if "City0" in t]
    assert hits, "subtitle text box did not render"
    subtitle = hits[0]
    rendered = [f"City{i}" in subtitle for i in range(7)]
    assert rendered[:3] == [True, True, True], "3-location cap is still a no-op"
    assert rendered[3:] == [False] * 4, "extra locations leaked through"


def test_ppt_locations_short_list_unchanged():
    report = _intel_report(locations=["Dallas, TX", "Austin, TX"])
    pptx_bytes = market_intel_reports.generate_intel_ppt(report)
    hits = [t for t in _pptx_texts(pptx_bytes) if "Dallas, TX" in t]
    assert hits and "Austin, TX" in hits[0]


# ===========================================================================
# Site 8: research.get_location_boards international "boards" cap
# ===========================================================================
def test_location_boards_capped_at_five_for_germany():
    """Germany's global_supply.json entry configures 12 boards -- repro
    before the fix showed all 12 "(CPC)" suffixed entries flowing through
    unbounded. After the fix, at most 5 come from that source."""
    boards = research.get_location_boards(["Berlin, Germany"])
    cpc_entries = [b for b in boards if b.endswith("(CPC)")]
    assert len(cpc_entries) == 5, (
        f"5-board cap is still a no-op (got {len(cpc_entries)} CPC entries): "
        f"{cpc_entries}"
    )
    # Sanity: the country-specific job-board data source really does have
    # more than 5 configured (otherwise this test wouldn't distinguish
    # capped from uncapped behavior).
    import json as _json

    _gs = _json.loads((PROJECT_ROOT / "data" / "global_supply.json").read_text())
    assert len(_gs["country_job_boards"]["Germany"]["boards"]) > 5


def test_location_boards_at_cap_boundary_unchanged_for_kenya():
    """Kenya's global_supply.json entry configures exactly 5 boards -- at
    the cap boundary, so slicing must be a true no-op here (proves the fix
    doesn't drop legitimate entries when already within budget)."""
    boards = research.get_location_boards(["Nairobi, Kenya"])
    cpc_entries = [b for b in boards if b.endswith("(CPC)")]
    assert len(cpc_entries) == 5

    import json as _json

    _gs = _json.loads((PROJECT_ROOT / "data" / "global_supply.json").read_text())
    assert len(_gs["country_job_boards"]["Kenya"]["boards"]) == 5


def test_location_boards_missing_country_entry_no_crash():
    """Norway is a recognized country (COUNTRY_DATA) with no entry in
    global_supply.json's country_job_boards -- the `.get("boards") or []`
    fallback must degrade gracefully rather than raising."""
    boards = research.get_location_boards(["Oslo, Norway"])
    assert isinstance(boards, list)
    assert any("Norway" in b for b in boards)  # LinkedIn/Indeed still appended
    assert not any(b.endswith("(CPC)") for b in boards)


if __name__ == "__main__":
    _failures = 0
    for _name, _fn in sorted(globals().items()):
        if _name.startswith("test_") and callable(_fn):
            try:
                _fn()
                print(f"PASS {_name}")
            except AssertionError as exc:
                _failures += 1
                print(f"FAIL {_name}: {exc}")
            except Exception as exc:  # noqa: BLE001
                _failures += 1
                print(f"ERROR {_name}: {exc}")
    if _failures:
        sys.exit(1)
