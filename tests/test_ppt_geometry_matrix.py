"""Structural-robustness matrix for ppt_generator.py (fix/gate-confidence-layout).

CONTEXT: a previous wave fixed three specific text collisions (slide 5
benchmark rows, slide 5 category card, slide 7 competitor Why/Counter --
see test_slide5_slide7_layout_and_claims.py) discovered on ONE client's
bundle. Those three fixes are real and must stay green, but they were found
one bundle at a time -- exactly the failure mode this file exists to close.
Every slide builder in ppt_generator.py positions content with geometry
(row pitches, card heights, item caps) that must hold across the FULL
realistic range of future plans, not just the bundles that happened to
surface a defect so far.

This file generates decks across that range -- channel counts (low/high),
location counts (1/many), competitor counts (0/1/3/more-than-cap), client/
role/industry-label name lengths (very long/very short), enrichment depth
(minimal/full), and wide-glyph/long-formatted currencies -- and asserts
real, independently-measured geometry on each:
  - nothing renders off the 13.333in x 7.5in canvas
  - no run is sized below the 8pt readability floor
  - no client-facing text reads as truncated mid-word (a letter glued
    directly to a trailing "..."/"…" with no separating space)
  - the specific slides this wave hardened (Cover, Push Meets Pull, Role
    Breakdown, Plan Comparison & Implementation, Quality Outcomes,
    Competitive Landscape) hold their measured-content-vs-declared-box and
    inter-element clearance invariants across the sparse/dense extremes
    that would have broken their OLD hardcoded constants.

A fully generic "do any two declared text-box rectangles overlap" check was
tried while writing this file and produced only false positives: several
boxes in this codebase are deliberately oversized CONTAINERS (e.g. the
Push/Pull detail textbox spans the card's full remaining height while the
real text only fills the first few lines, with the sibling "This plan:"
split-line box correctly positioned from the MEASURED line count, not the
container's declared height) -- checking declared boxes flags these as
"overlapping" when the real rendered text never does. So overlap/overflow
here is checked the same way test_slide5_slide7_layout_and_claims.py
already does it: using each shape's OWN measured content extent against
the specific neighbour it could plausibly collide with, not a blanket
box-vs-box scan.

Runs under pytest, or standalone: ``python3 tests/test_ppt_geometry_matrix.py``.
"""

from __future__ import annotations

import io
import re
import sys
from pathlib import Path
from typing import Any, Callable, Dict

import pytest

PROJECT_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(PROJECT_ROOT))

import ppt_generator as ppt  # noqa: E402
import tools_regen_bundles as T  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE  # noqa: E402

EMU_PER_IN = 914400
TOL = 0.03
SLIDE_W_IN = ppt.SLIDE_WIDTH / EMU_PER_IN
SLIDE_H_IN = ppt.SLIDE_HEIGHT / EMU_PER_IN

_ALL_CHANNEL_KEYS = (
    "programmatic_dsp",
    "global_boards",
    "niche_boards",
    "social_media",
    "regional_boards",
    "employer_branding",
    "apac_regional",
    "emea_regional",
)

# Same regex bundle_qa.py's own mid-word-truncation gate uses (a letter
# directly, unspaced, against a trailing ellipsis) -- reproduced
# independently here (not imported) since this file owns only
# ppt_generator.py/tests, not bundle_qa.py.
_MID_WORD_ELLIPSIS_RE = re.compile(r"[a-zA-Z](\.\.\.|…)$")


# ---------------------------------------------------------------------------
# Shape helpers (same conventions as test_slide5_slide7_layout_and_claims.py)
# ---------------------------------------------------------------------------
def _iter_shapes(shapes):
    for sh in shapes:
        yield sh
        if sh.shape_type == 6:  # group
            try:
                yield from _iter_shapes(sh.shapes)
            except Exception:
                pass


def _text_shapes(slide):
    out = []
    for sh in _iter_shapes(slide.shapes):
        try:
            if sh.has_text_frame and sh.text_frame.text.strip():
                out.append(sh)
        except Exception:
            pass
    return out


def _slide_by_headline(prs: Presentation, headline: str):
    for slide in prs.slides:
        for sh in _iter_shapes(slide.shapes):
            if not (sh.has_text_frame and sh.text_frame.text.strip()):
                continue
            if sh.text_frame.text.strip() != headline:
                continue
            top_in = sh.top / EMU_PER_IN
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.size and abs(r.font.size.pt - 26.0) < 0.1 and 0.4 < top_in < 0.5:
                        return slide
    return None


def _run_font_pt(shape, default: float = 10.0) -> float:
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if r.font.size is not None:
                return r.font.size.pt
    return default


def _content_h_in(text: str, width_in: float, font_pt: float, char_em: float = 0.53) -> float:
    """Same independent estimate test_slide5_slide7_layout_and_claims.py
    uses: ppt_generator's own ``_estimate_lines`` against the ACTUAL
    generated font size/width, times the calibrated 1.42 line-height factor
    ``_autofit_textframe`` documents."""
    n_lines = ppt._estimate_lines(text, width_in, font_pt, char_em=char_em)
    return n_lines * (font_pt * 1.42) / 72.0


def _rounded_rect_cards(slide, min_w_in: float, min_h_in: float):
    cards = []
    for sh in _iter_shapes(slide.shapes):
        if sh.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            continue
        try:
            if sh.auto_shape_type != MSO_SHAPE.ROUNDED_RECTANGLE:
                continue
        except Exception:
            continue
        if sh.has_text_frame and sh.text_frame.text.strip():
            continue
        if sh.width / EMU_PER_IN < min_w_in or sh.height / EMU_PER_IN < min_h_in:
            continue
        cards.append(sh)
    cards.sort(key=lambda s: s.top)
    return cards


# ---------------------------------------------------------------------------
# Matrix fixtures
# ---------------------------------------------------------------------------
def _base_plan(**overrides) -> Dict[str, Any]:
    d: Dict[str, Any] = {
        "client_name": "Acme",
        "industry": "hospitality_travel",
        "industry_label": "Hospitality & Travel",
        "locations": ["London, United Kingdom"],
        "roles": ["Commercial Cab Driver"],
        "budget": "£2,000,000",
        "work_environment": "hybrid",
        "channel_categories": {"programmatic_dsp": True, "global_boards": True, "social_media": True},
    }
    d.update(overrides)
    return d


def _pipeline_plan(brief: Dict[str, Any]) -> Dict[str, Any]:
    """Runs the SAME app.py -> budget_engine -> gold_standard pipeline
    tools_regen_bundles.py uses, for full-enrichment coverage (a real
    _budget_allocation + _gold_standard, not a bare dict)."""
    brief = dict(brief)
    brief.setdefault("requester_name", "Test Runner")
    brief.setdefault("requester_email", "test@example.com")
    return T.build_plan_data(brief)


def _roles(n: int, title: str = "Warehouse Associate"):
    return [
        {"title": f"{title} {i + 1}" if n > 1 else title, "count": 50, "tier": "Hourly"}
        for i in range(n)
    ]


MATRIX: Dict[str, Callable[[], Dict[str, Any]]] = {
    # -- channel counts: low / high --
    "two_channels": lambda: _base_plan(
        channel_categories={"programmatic_dsp": True, "global_boards": True}
    ),
    "all_channels": lambda: _base_plan(
        channel_categories={k: True for k in _ALL_CHANNEL_KEYS}
    ),
    # -- locations: 1 / many --
    "one_location": lambda: _base_plan(locations=["London, United Kingdom"]),
    "many_locations": lambda: _base_plan(
        locations=[f"City{i}, Country{i}" for i in range(9)]
    ),
    # -- competitors: 0 / 1 / 3 (= cap) / more than cap --
    "zero_competitors": lambda: _base_plan(competitors=[]),
    "one_competitor": lambda: _base_plan(competitors=["Marriott"]),
    "three_competitors_long_desc": lambda: _base_plan(
        competitors=[
            {
                "name": "Marriott International Hospitality Group",
                "description": (
                    "A very long raw description with no dash or semicolon "
                    "separators whatsoever just prose that keeps going and "
                    "going without any punctuation marks that would count "
                    "as clause boundaries for the truncation helper to "
                    "latch onto at all here we go more words"
                ),
            },
            {
                "name": "Hilton Worldwide Holdings Incorporated",
                "description": (
                    "Another long raw description with no dash or semicolon "
                    "separators whatsoever just prose that keeps going and "
                    "going without any punctuation marks that would count "
                    "as clause boundaries for the truncation helper to "
                    "latch onto at all here we go more words"
                ),
            },
            {
                "name": "Hyatt Hotels Corporation International",
                "description": (
                    "A third long raw description with no dash or semicolon "
                    "separators whatsoever just prose that keeps going and "
                    "going without any punctuation marks that would count "
                    "as clause boundaries for the truncation helper to "
                    "latch onto at all here we go more words"
                ),
            },
        ],
    ),
    "more_than_cap_competitors": lambda: _base_plan(
        competitors=[
            "Marriott International Hospitality Holdings",
            "Hilton Worldwide Holdings Incorporated",
            "Hyatt Hotels Corporation International",
            "IHG Hotels and Resorts Group PLC",
            "Accor Group SA",
            "Wyndham Hotels and Resorts",
        ]
    ),
    # -- name lengths: very long / very short (client, role, industry label) --
    "very_long_client_name": lambda: _base_plan(
        client_name=(
            "The International Consolidated Hospitality Leisure and Gaming "
            "Holdings Group PLC & Co KGaA"
        )
    ),
    "very_short_client_name": lambda: _base_plan(client_name="Zo"),
    "very_long_role_title": lambda: _base_plan(
        roles=[
            "Senior Regional Assistant Vice President of Front-of-House "
            "Guest Experience Operations"
        ]
    ),
    "very_short_role_title": lambda: _base_plan(roles=["Aide"]),
    "long_industry_label": lambda: _base_plan(
        industry_label="Hospitality, Travel, Leisure, Gaming and Entertainment Services (Global)"
    ),
    # -- enrichment depth: minimal / full --
    "minimal_enrichment": lambda: {
        "client_name": "MiniCo",
        "industry": "general_entry_level",
        "locations": ["Austin, TX"],
        "roles": ["Clerk"],
        "budget": "$5,000",
    },
    "full_enrichment_dense": lambda: _pipeline_plan(
        dict(
            client_name="Amazon",
            industry="Logistics & Supply Chain",
            budget="$10,000,000",
            campaign_duration="1 year",
            hire_volume="5000+ hires",
            work_environment="onsite",
            locations=[
                "Seattle, WA", "Denver, CO", "Newark, NJ", "Memphis, TN",
                "Dallas, TX", "Atlanta, GA", "Phoenix, AZ", "Chicago, IL",
            ],
            roles=[r["title"] for r in _roles(14)],
            target_roles=_roles(14),
            channel_categories={k: True for k in _ALL_CHANNEL_KEYS},
            competitors=["FedEx", "UPS", "XPO Logistics", "DHL", "Ryder", "J.B. Hunt"],
        )
    ),
    # Render-verified (Keynote) reproduction: a very long client name
    # interpolated into the "Campaign setup & publisher activation for
    # {client}" timeline bullet, COMBINED with enough real budget-allocation
    # data to fill the comparison panel to its row cap -- the exact
    # precondition that pushed the phase card's bullet text past the card
    # and into the footer text on origin/main.
    "full_enrichment_long_name_dense": lambda: _pipeline_plan(
        dict(
            client_name=(
                "The International Consolidated Logistics and Supply Chain "
                "Holdings Group PLC"
            ),
            industry="Logistics & Supply Chain",
            budget="$10,000,000",
            campaign_duration="1 year",
            hire_volume="5000+ hires",
            work_environment="onsite",
            locations=[
                "Seattle, WA", "Denver, CO", "Newark, NJ", "Memphis, TN",
                "Dallas, TX", "Atlanta, GA", "Phoenix, AZ", "Chicago, IL",
            ],
            roles=[r["title"] for r in _roles(14)],
            target_roles=_roles(14),
            channel_categories={k: True for k in _ALL_CHANNEL_KEYS},
            competitors=["FedEx", "UPS", "XPO Logistics", "DHL", "Ryder", "J.B. Hunt"],
        )
    ),
    "full_enrichment_goal_gap": lambda: _pipeline_plan(
        dict(
            client_name="Mercy Health",
            industry="Healthcare & Medical",
            budget="$250,000",
            campaign_duration="18 months",
            hire_volume="200+ hires",
            work_environment="onsite",
            locations=["Chicago, IL"],
            roles=["Registered Nurse"],
            target_roles=[{"title": "Registered Nurse", "count": 200, "tier": "Salaried"}],
        )
    ),
    "full_enrichment_10_roles": lambda: _pipeline_plan(
        dict(
            client_name="atria Senior living",
            industry="Healthcare & Medical",
            budget="$300,000",
            campaign_duration="18 months",
            hire_volume="500+ hires",
            work_environment="remote",
            locations=["New York, NY"],
            roles=[
                "Memory Care Associate", "Nurse", "Cook", "Driver",
                "Maintenance Technician", "Server/Waitstaff", "Shift/Charge Nurse",
                "Dishwasher", "Housekeeper", "Sales",
            ],
            target_roles=[
                {"title": t, "count": 50, "tier": "Hourly"}
                for t in (
                    "Memory Care Associate", "Nurse", "Cook", "Driver",
                    "Maintenance Technician", "Server/Waitstaff", "Shift/Charge Nurse",
                    "Dishwasher", "Housekeeper", "Sales",
                )
            ],
            competitors=["Brookdale Senior Living", "Sunrise Senior Living", "Amazon"],
        )
    ),
    # -- wide-glyph currency / long formatted numbers --
    "wide_currency_idr": lambda: _pipeline_plan(
        dict(
            client_name="PT Gojek Indonesia",
            industry="Technology & Engineering",
            budget="Rp15,000,000,000",
            campaign_duration="6 months",
            hire_volume="300+ hires",
            work_environment="onsite",
            locations=["Jakarta"],
            roles=["Driver Partner"],
            target_roles=[{"title": "Driver Partner", "count": 300, "tier": "Hourly"}],
        )
    ),
    "gbp_narrow_budget": lambda: _base_plan(budget="£2,000,000"),
}


@pytest.fixture(scope="module")
def decks() -> Dict[str, Presentation]:
    """Generate every matrix shape's deck once (module-scoped: this is a
    render-and-measure suite, not a unit test -- re-generating ~17 full
    decks per test would multiply the real cost of pptx serialization for
    no additional coverage)."""
    out = {}
    for name, build in MATRIX.items():
        data = build()
        pptx_bytes = ppt.generate_pptx(data)
        out[name] = Presentation(io.BytesIO(pptx_bytes))
    return out


# ---------------------------------------------------------------------------
# Universal checks -- safe across every shape/slide (no line-height modelling
# assumptions, so no false-positive risk from oversized-container textboxes).
# ---------------------------------------------------------------------------
class TestUniversalGeometry:
    def test_no_shape_off_canvas(self, decks):
        violations = []
        for name, prs in decks.items():
            for si, slide in enumerate(prs.slides):
                for sh in _iter_shapes(slide.shapes):
                    try:
                        l, t, w, h = sh.left, sh.top, sh.width, sh.height
                    except Exception:
                        continue
                    if l is None or t is None or w is None or h is None:
                        continue
                    l_in, t_in = l / EMU_PER_IN, t / EMU_PER_IN
                    r_in, b_in = (l + w) / EMU_PER_IN, (t + h) / EMU_PER_IN
                    if (
                        l_in < -TOL
                        or t_in < -TOL
                        or r_in > SLIDE_W_IN + TOL
                        or b_in > SLIDE_H_IN + TOL
                    ):
                        violations.append(
                            f"{name} slide {si + 1}: shape at l={l_in:.2f} t={t_in:.2f} "
                            f"r={r_in:.2f} b={b_in:.2f} exceeds canvas "
                            f"({SLIDE_W_IN:.2f}x{SLIDE_H_IN:.2f}in)"
                        )
        assert not violations, "\n".join(violations)

    def test_no_sub_8pt_runs(self, decks):
        violations = []
        for name, prs in decks.items():
            for si, slide in enumerate(prs.slides):
                for sh in _text_shapes(slide):
                    for p in sh.text_frame.paragraphs:
                        for r in p.runs:
                            if r.font.size is not None and r.font.size.pt < 7.9:
                                violations.append(
                                    f"{name} slide {si + 1}: {r.font.size.pt}pt run "
                                    f"{r.text!r}"
                                )
        assert not violations, "\n".join(violations)

    def test_no_mid_word_ellipsis_anywhere(self, decks):
        violations = []
        for name, prs in decks.items():
            for si, slide in enumerate(prs.slides):
                for sh in _text_shapes(slide):
                    t = sh.text_frame.text.strip()
                    if _MID_WORD_ELLIPSIS_RE.search(t):
                        violations.append(f"{name} slide {si + 1}: {t[-60:]!r}")
        assert not violations, "\n".join(violations)


# ---------------------------------------------------------------------------
# Targeted checks -- the specific slides this wave hardened. Each mirrors
# test_slide5_slide7_layout_and_claims.py's pattern: measure the shape's
# OWN actual content against the ACTUAL declared geometry of a specific
# neighbour, not a blanket box-vs-box scan.
# ---------------------------------------------------------------------------
class TestCoverSlideClientNameFit:
    def test_long_client_name_does_not_collide_with_industry_line(self, decks):
        prs = decks["very_long_client_name"]
        slide = prs.slides[0]
        client_sh = next(
            (
                sh
                for sh in _text_shapes(slide)
                if sh.text_frame.text.strip().startswith("The International Consolidated")
            ),
            None,
        )
        industry_sh = next(
            (
                sh
                for sh in _text_shapes(slide)
                if sh.text_frame.text.strip() == "Hospitality & Travel"
            ),
            None,
        )
        assert client_sh is not None, "expected the client hero text on the cover slide"
        assert industry_sh is not None, "expected the industry subtitle on the cover slide"
        font_pt = _run_font_pt(client_sh)
        width_in = client_sh.width / EMU_PER_IN
        needed_in = _content_h_in(client_sh.text_frame.text, width_in, font_pt, char_em=0.53)
        client_bottom_in = client_sh.top / EMU_PER_IN + needed_in
        industry_top_in = industry_sh.top / EMU_PER_IN
        assert client_bottom_in <= industry_top_in + TOL, (
            f"client hero text needs to reach {client_bottom_in:.2f}in but the "
            f"industry subtitle starts at {industry_top_in:.2f}in"
        )

    def test_short_client_name_renders_at_original_anchor(self, decks):
        """False-positive guard: a normal-length name must NOT be shrunk or
        shifted -- the fit-to-2-lines/cascade logic only engages when
        actually needed."""
        prs = decks["very_short_client_name"]
        slide = prs.slides[0]
        client_sh = next(
            (sh for sh in _text_shapes(slide) if sh.text_frame.text.strip() == "Zo"), None
        )
        assert client_sh is not None
        assert _run_font_pt(client_sh) == pytest.approx(42.0)
        assert client_sh.top / EMU_PER_IN == pytest.approx(3.48, abs=0.01)


class TestPushMeetsPullEnvelope:
    def test_card_never_exceeds_ceiling_even_with_every_channel_itemized(self, decks):
        """_push_pull_split_line's contract is "every channel MUST be
        itemized" -- the all_channels shape maximizes both push and pull
        itemized-list length. The card must grow to fit it WITHOUT the
        declared card height exceeding the old 4.6in ceiling."""
        prs = decks["all_channels"]
        slide = _slide_by_headline(prs, "Push Meets Pull")
        assert slide is not None
        cards = _rounded_rect_cards(slide, min_w_in=5.0, min_h_in=1.0)
        assert cards, "expected at least one Push/Pull card"
        for c in cards:
            assert c.height / EMU_PER_IN <= 4.6 + TOL, (
                f"Push/Pull card height {c.height / EMU_PER_IN:.2f}in exceeds the 4.6in ceiling"
            )

    def test_single_populated_side_widens_instead_of_stranding_half_the_slide(self):
        """Defensive case: a future deck-KB edit (or caller) supplying only
        ONE of push/pull must not draw a single 6in card at its normal
        fixed X, leaving the other half of the slide blank."""
        prs = Presentation()
        prs.slide_width = ppt.SLIDE_WIDTH
        prs.slide_height = ppt.SLIDE_HEIGHT
        deck_kb = dict(ppt._load_deck_kb())
        pmp = dict(deck_kb.get("push_meets_pull") or {})
        pmp["pull"] = {}  # simulate the pull section being absent
        deck_kb["push_meets_pull"] = pmp
        data = _base_plan()
        ppt._build_slide_push_meets_pull(prs, data, deck_kb)
        assert len(prs.slides) == 1
        slide = prs.slides[0]
        cards = _rounded_rect_cards(slide, min_w_in=5.0, min_h_in=1.0)
        assert len(cards) == 1, "expected exactly one card when only one side has content"
        card = cards[0]
        card_w_in = card.width / EMU_PER_IN
        assert card_w_in > 6.5, (
            f"single populated side rendered at {card_w_in:.2f}in wide -- "
            "still using the two-column 6in width, stranding the other half blank"
        )


class TestRoleBreakdownEnvelope:
    def test_many_roles_never_run_past_the_footer(self, decks):
        prs = decks["full_enrichment_dense"]  # 14 roles -- past the old rows[:12] cap
        slide = _slide_by_headline(prs, "Role Breakdown")
        if slide is None:
            pytest.skip("Role Breakdown did not render for this fixture (CPA Reference took the slot)")
        table_bottoms = [
            sh.top / EMU_PER_IN + sh.height / EMU_PER_IN
            for sh in _iter_shapes(slide.shapes)
            if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
            and not (sh.has_text_frame and sh.text_frame.text.strip())
            and sh.height / EMU_PER_IN < 1.0  # exclude the full-slide background rect
            and 11.9 < sh.width / EMU_PER_IN < 12.1  # table-row-width rects only (excludes the footer's own 12.23in-wide rule)
        ]
        assert table_bottoms, "expected filled row/header rectangles on the Role Breakdown table"
        assert max(table_bottoms) <= 7.12 + TOL, (
            f"Role Breakdown table row extends to {max(table_bottoms):.2f}in, "
            "past the footer rule at 7.12in"
        )
        # A rollup ("+N more roles") must appear when roles were dropped.
        texts = " ".join(sh.text_frame.text for sh in _text_shapes(slide))
        if "more role" in texts:
            assert "workbook" in texts.lower()

    def test_minimum_eligible_role_count_is_vertically_centered_not_top_stranded(self, decks):
        prs = decks["full_enrichment_10_roles"]
        slide = _slide_by_headline(prs, "Role Breakdown")
        assert slide is not None
        header_bgs = [
            sh
            for sh in _iter_shapes(slide.shapes)
            if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
            and not (sh.has_text_frame and sh.text_frame.text.strip())
            and sh.height / EMU_PER_IN < 1.0  # exclude the full-slide background rect
            and 11.9 < sh.width / EMU_PER_IN < 12.1  # table-row-width rects only (excludes the footer's own 12.23in-wide rule)
        ]
        assert header_bgs
        table_top_in = min(sh.top / EMU_PER_IN for sh in header_bgs)
        table_bottom_in = max(sh.top / EMU_PER_IN + sh.height / EMU_PER_IN for sh in header_bgs)
        # never above the old fixed anchor
        assert table_top_in >= 1.7 - TOL
        # 10 rows is dense enough it should sit close to the available
        # band, not stranded with a huge gap either above or below
        assert table_bottom_in <= 6.95 + TOL


class TestComparisonTimelineEnvelope:
    def test_goal_band_clears_the_rounded_panels_corner_radius(self, decks):
        """The comparison panel used to be a flat fixed 2.95in regardless of
        row count + the optional goal band. The panel is a ROUNDED_RECTANGLE
        auto-shape with python-pptx's default corner radius (~16.667% of its
        shorter side -- ~0.49in at the old fixed 2.95in height); the flat,
        square-cornered amber "CLIENT GOAL" band spans the FULL WIDTH of
        both panels, so whenever its bottom edge sits close to the panel's
        own bottom edge, its square corners sit where the panel's rounded
        corner has already curved inward -- visually, the panel's own white
        rounded corner peeks out from beneath the band (confirmed via a
        real Keynote render of this exact atria-shaped bundle). This was
        true for EVERY goal-gap plan on origin/main, not just an overflow
        edge case: comparison_rows is always capped at exactly 4 whenever a
        goal row exists, so the old fixed panel always left only ~0.2in of
        clearance below the band -- well inside the corner-radius zone.
        The fix's content-derived panel height adds real bottom clearance
        (~0.6in here) that clears it; this asserts a minimum clearance
        rather than mere non-overflow, since non-overflow alone doesn't
        capture the corner-radius visual defect (it was already
        "non-overflowing" on origin/main by the bounding-box test alone).
        """
        # full_enrichment_10_roles (the atria-shaped fixture) is the one
        # that reaches 4 comparison rows + a goal band (500-hire goal vs. a
        # much lower projection) -- the real precondition this bug needs.
        prs = decks["full_enrichment_10_roles"]
        slide = _slide_by_headline(prs, "Plan Comparison & Implementation")
        assert slide is not None
        goal_band = next(
            (
                sh
                for sh in _text_shapes(slide)
                if "CLIENT GOAL" in sh.text_frame.text
            ),
            None,
        )
        assert goal_band is not None, (
            "expected this plan's real hire-goal shortfall to render a "
            "CLIENT GOAL band -- fixture no longer reproduces the goal-gap "
            "precondition this test needs"
        )
        panels = _rounded_rect_cards(slide, min_w_in=5.0, min_h_in=1.5)
        assert panels, "expected the two comparison panels"
        goal_band_bottom_in = goal_band.top / EMU_PER_IN + goal_band.height / EMU_PER_IN
        for panel in panels:
            panel_bottom_in = panel.top / EMU_PER_IN + panel.height / EMU_PER_IN
            clearance_in = panel_bottom_in - goal_band_bottom_in
            assert clearance_in >= 0.35, (
                f"CLIENT GOAL band bottom {goal_band_bottom_in:.2f}in leaves only "
                f"{clearance_in:.2f}in of clearance before panel bottom "
                f"{panel_bottom_in:.2f}in -- inside the rounded panel's own "
                "~0.49in corner-radius curvature zone"
            )

    def test_legend_never_overlaps_the_implementation_timeline_header(self, decks):
        for name in ("full_enrichment_goal_gap", "full_enrichment_dense", "one_location"):
            prs = decks[name]
            slide = _slide_by_headline(prs, "Plan Comparison & Implementation")
            assert slide is not None
            legend = next(
                (sh for sh in _text_shapes(slide) if "Beating benchmark" in sh.text_frame.text),
                None,
            )
            header = next(
                (
                    sh
                    for sh in _text_shapes(slide)
                    if sh.text_frame.text.strip() == "IMPLEMENTATION TIMELINE"
                ),
                None,
            )
            assert legend is not None and header is not None
            legend_bottom_in = legend.top / EMU_PER_IN + legend.height / EMU_PER_IN
            header_top_in = header.top / EMU_PER_IN
            assert legend_bottom_in <= header_top_in + TOL, (
                f"[{name}] legend bottom {legend_bottom_in:.2f}in overlaps "
                f"IMPLEMENTATION TIMELINE header at {header_top_in:.2f}in"
            )

    def test_phase_cards_never_run_into_the_footer(self, decks):
        """_interpolate_timeline_bullets splices the client name/top channel
        labels into these bullets -- a long one wraps to 2 lines, and (on a
        plan whose comparison panel is also tall) used to push the phase
        card's bottom edge past the footer rule."""
        for name, prs in decks.items():
            slide = _slide_by_headline(prs, "Plan Comparison & Implementation")
            if slide is None:
                continue
            phase_cards = _rounded_rect_cards(slide, min_w_in=3.0, min_h_in=1.0)
            for card in phase_cards:
                bottom_in = card.top / EMU_PER_IN + card.height / EMU_PER_IN
                assert bottom_in <= 7.12 + TOL, (
                    f"[{name}] phase card bottom {bottom_in:.2f}in runs into the "
                    "footer rule at 7.12in"
                )

    def test_phase_bullet_text_never_overflows_its_own_box(self, decks):
        """PowerPoint textboxes do not clip -- an oversized bullet list
        renders PAST its declared box and (render-verified via Keynote on
        the very-long-client-name + full-enrichment combination) past the
        phase card itself and into the footer text. This measures the
        ACTUAL bullet content against its OWN declared box height, the
        same invariant Task 3 requires everywhere else."""
        violations = []
        for name, prs in decks.items():
            slide = _slide_by_headline(prs, "Plan Comparison & Implementation")
            if slide is None:
                continue
            for sh in _text_shapes(slide):
                if not sh.text_frame.text.strip().startswith("✓"):
                    continue
                width_in = sh.width / EMU_PER_IN
                # Multi-paragraph bullet list -- estimate each line's own
                # font size/text and sum, matching _autofit_textframe's own
                # per-paragraph measurement model.
                needed_in = 0.10  # textbox's own default top+bottom inset (0.05in each)
                for p in sh.text_frame.paragraphs:
                    runs = [r for r in p.runs if r.text]
                    if not runs:
                        continue
                    font_pt = max(r.font.size.pt if r.font.size else 8 for r in runs)
                    ptext = "".join(r.text for r in runs)
                    n = ppt._estimate_lines(ptext, width_in, font_pt)
                    needed_in += n * (font_pt * 1.42) / 72.0
                    sb = (p.space_before.pt if p.space_before else 0) / 72.0
                    sa = (p.space_after.pt if p.space_after else 0) / 72.0
                    needed_in += sb + sa
                height_in = sh.height / EMU_PER_IN
                if needed_in > height_in + TOL:
                    violations.append(
                        f"[{name}] phase bullet box needs ~{needed_in:.2f}in but is "
                        f"only {height_in:.2f}in tall: {sh.text_frame.text!r}"
                    )
        assert not violations, "\n".join(violations)


class TestQualityOutcomesRollup:
    def test_more_than_five_channels_get_an_explicit_rollup_not_silent_drop(self, decks):
        prs = decks["all_channels"]  # 8 channels enabled, no _budget_allocation
        slide = _slide_by_headline(prs, "Quality & ROI Projections")
        assert slide is not None, "expected the Quality Outcomes fallback slide"
        texts = " ".join(sh.text_frame.text for sh in _text_shapes(slide))
        assert "smaller channels" in texts, (
            "8 enabled channels exceed the 5-row cap -- expected an explicit "
            "'+N smaller channels' rollup row instead of a silent drop"
        )
        # 8 enabled channels, _QO_MAX_VISIBLE_ROWS=5 -> 4 visible + 1 rollup
        # row absorbing the remaining 4.
        m = re.search(r"\+(\d+) smaller channels", texts)
        assert m and int(m.group(1)) == 4, f"expected '+4 smaller channels', got: {texts!r}"


class TestCompetitiveLandscapeEnvelope:
    def test_competitor_card_stack_clears_source_line_and_footer(self, decks):
        """3 real competitor cards with longer counter-strategy prose put
        the last card's own bottom at ~6.24in -- the tallest content this
        slide renders. Reproduced directly against origin/main. The
        (now-retired) Market Positioning band used to be checked here too;
        with it gone, the Source line is anchored straight off this same
        measured card-stack bottom (see ppt_generator.py's
        retire(dead-positioning-band) comment), so both must still clear
        the footer and never collide with each other."""
        prs = decks["three_competitors_long_desc"]
        slide = _slide_by_headline(prs, "Competitive Landscape")
        assert slide is not None
        cards = _rounded_rect_cards(slide, min_w_in=4.0, min_h_in=1.0)
        assert cards, "expected competitor cards"
        cards_bottom_in = max(c.top / EMU_PER_IN + c.height / EMU_PER_IN for c in cards)
        # Card stack itself must stay clear of the Source line / footer.
        assert cards_bottom_in <= 7.12 + TOL

        source_line = next(
            (
                sh
                for sh in _text_shapes(slide)
                if sh.text_frame.text.strip().startswith("Sources:")
            ),
            None,
        )
        assert source_line is not None, "expected the Sources line"
        source_top_in = source_line.top / EMU_PER_IN
        assert source_top_in >= cards_bottom_in - TOL, (
            f"Sources line starts at {source_top_in:.2f}in but the "
            f"competitor card stack extends to {cards_bottom_in:.2f}in"
        )

    def test_competitor_count_variants_all_stay_on_canvas_and_uncollided(self, decks):
        for name in ("zero_competitors", "one_competitor", "more_than_cap_competitors"):
            prs = decks[name]
            slide = _slide_by_headline(prs, "Competitive Landscape")
            assert slide is not None, f"[{name}] expected the Competitive Landscape slide"
            cards = _rounded_rect_cards(slide, min_w_in=4.0, min_h_in=1.0)
            # more_than_cap_competitors must still cap at _MAX_COMPETITOR_CARDS (3)
            if name == "more_than_cap_competitors":
                assert len(cards) <= 3
            for a, b in zip(cards, cards[1:]):
                a_bottom_in = a.top / EMU_PER_IN + a.height / EMU_PER_IN
                b_top_in = b.top / EMU_PER_IN
                assert a_bottom_in <= b_top_in + TOL, (
                    f"[{name}] competitor cards overlap: {a_bottom_in:.2f}in vs {b_top_in:.2f}in"
                )


# ---------------------------------------------------------------------------
# Regression guard: the three already-fixed collisions stay fixed under the
# WIDER matrix too (not just the one Uber-shaped fixture the original fix
# was verified against).
# ---------------------------------------------------------------------------
class TestExistingCollisionFixesHoldAcrossMatrix:
    def test_slide5_benchmark_and_category_rows_never_overlap_across_matrix(self, decks):
        for name, prs in decks.items():
            slide = _slide_by_headline(prs, "Channel Strategy & Investment")
            if slide is None:
                continue
            value_cells = [
                sh
                for sh in _text_shapes(slide)
                if 9.7 < sh.left / EMU_PER_IN < 10.1 and 1.9 < sh.top / EMU_PER_IN < 5.0
            ]
            value_cells.sort(key=lambda s: s.top)
            for a, b in zip(value_cells, value_cells[1:]):
                a_bottom_in = a.top / EMU_PER_IN + a.height / EMU_PER_IN
                b_top_in = b.top / EMU_PER_IN
                assert a_bottom_in <= b_top_in + TOL, (
                    f"[{name}] benchmark row overlap at {a_bottom_in:.2f}in vs {b_top_in:.2f}in"
                )

    def test_slide7_why_counter_never_overlap_across_matrix(self, decks):
        for name, prs in decks.items():
            slide = _slide_by_headline(prs, "Competitive Landscape")
            if slide is None:
                continue
            why_shapes = sorted(
                (sh for sh in _text_shapes(slide) if sh.text_frame.text.startswith("Why:")),
                key=lambda s: s.top,
            )
            counter_shapes = sorted(
                (sh for sh in _text_shapes(slide) if sh.text_frame.text.startswith("Counter:")),
                key=lambda s: s.top,
            )
            for why_sh, counter_sh in zip(why_shapes, counter_shapes):
                width_in = why_sh.width / EMU_PER_IN
                font_pt = _run_font_pt(why_sh)
                why_needed_in = _content_h_in(why_sh.text_frame.text, width_in, font_pt)
                why_bottom_in = why_sh.top / EMU_PER_IN + why_needed_in
                counter_top_in = counter_sh.top / EMU_PER_IN
                assert why_bottom_in <= counter_top_in + TOL, (
                    f"[{name}] Why text overprints Counter: needs {why_bottom_in:.2f}in, "
                    f"Counter starts at {counter_top_in:.2f}in"
                )


if __name__ == "__main__":
    raise SystemExit(pytest.main([__file__, "-v"]))
