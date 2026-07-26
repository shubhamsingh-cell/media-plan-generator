"""Regression tests for 3 of the 4 currency-integrity defects found by
regenerating the real Uber bundle (GBP 2,000,000 plan, 6 non-US markets:
UK/Australia/Mexico/Argentina/Canada/New Zealand) and running
``bundle_qa.run_bundle_qa`` on the output. That run reported 13
``currency_symbol_mixing`` criticals plus one on deck slide 2:

  Defect 1 -- Market Intelligence!F17:F22 (the "Location Intelligence"
  table's Median Income column) rendered a bare "$" for a value that is
  actually each market's OWN local-currency median salary (research.
  COUNTRY_DATA, e.g. United Kingdom 42000/"GBP"). The sibling "Location
  Economic Context" table on the SAME sheet had already been fixed to
  render "£42,000 (GBP)" via plan_currency.symbol_for_code -- this column
  was missed.

  Defect 2 -- Quality Intelligence!H8:H13 (the "City-Level Supply-Demand
  Data" table's Salary Range column) had the identical bug: excel_v2.
  _salary_range_from_per_role hardcoded a literal "$" on both ends of the
  range regardless of which market it described.

  Defect 3 -- deck slide 2's COMPLICATION card renders a genuinely
  USD-sourced KB benchmark string ("CDL/last-mile roles most expensive
  within the range, at $25-$50 CPA") verbatim, with no USD marker, right
  next to this plan's own "£" figures elsewhere on the slide -- implying a
  conversion that never happened. This string legitimately traces to a USD
  source (data/recruitment_benchmarks_deep.json) and must NOT be converted
  or relabelled to "£"; it must be marked honest ("US$25-US$50") the same
  way ppt_generator's slide 5 benchmark table already marks USD-sourced
  rows via ``_mark_usd`` / the ``_get_active_currency() != "USD"`` gate.

Defect 4 (the exec_summary_budget_footing_mismatch false positive on
bundle_qa's own parser) is covered separately in
tests/test_bundle_qa_money_suffix_parsing.py.

VACUOUSNESS: every test in this file was run against a pre-fix throwaway
worktree (``git worktree add --detach HEAD``) at the parent commit plus
every OTHER agent's uncommitted work on this tree, with only this file's
target fixes reverted -- see the task report for the observed pre-fix
failures.
"""

from __future__ import annotations

import io
import os
import sys

import pytest

PROJECT_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, PROJECT_ROOT)

import excel_v2  # noqa: E402
import ppt_generator  # noqa: E402
import research as research_mod  # noqa: E402
import tools_regen_bundles as T  # noqa: E402
from bundle_qa import run_bundle_qa  # noqa: E402
from pptx import Presentation  # noqa: E402
from openpyxl import load_workbook  # noqa: E402


# The real diagnosed brief: GBP budget, 6 non-US markets, rideshare driver role.
UBER_GBP_BRIEF: dict = {
    "client_name": "uber",
    "requester_name": "Test Requester",
    "requester_email": "test@joveo.com",
    "budget": "£2,000,000",
    "campaign_duration": "1-3 months",
    "hire_volume": "500+ hires",
    "work_environment": "hybrid",
    "locations": ["UK", "Australia", "Mexico", "argentina", "canada", "new zealand"],
    "roles": ["commercial cab driver"],
    "target_roles": [
        {"title": "commercial cab driver", "count": 500, "tier": "Hourly"}
    ],
    "notes": "Global commercial driver supply across 6 non-US markets.",
}

# Exact values verbatim from the diagnosed incident (research.COUNTRY_DATA's
# own median_salary/currency for each of these 6 countries).
_EXPECTED_MEDIAN_INCOME = {
    "UK": "£42,000 (GBP)",
    "Australia": "A$55,000 (AUD)",
    "Mexico": "MX$9,500 (MXN)",
    "argentina": "AR$8,000 (ARS)",
    "canada": "C$52,000 (CAD)",
    "new zealand": "NZ$48,000 (NZD)",
}

# Quality Intelligence's "Market" column is title-cased (excel_v2._title_case_city).
_EXPECTED_SALARY_RANGE = {
    "Uk": "£25,500 - £42,500 (GBP)",
    "Australia": "A$27,000 - A$45,000 (AUD)",
    "Mexico": "MX$6,000 - MX$10,000 (MXN)",
    "Argentina": "AR$3,600 - AR$6,000 (ARS)",
    "Canada": "C$25,500 - C$42,500 (CAD)",
    "New Zealand": "NZ$22,500 - NZ$37,500 (NZD)",
}


def _generate_bundle(brief: dict) -> dict:
    data = T.build_plan_data(brief)
    xlsx_obj = excel_v2.generate_excel_v2(data, research_mod=research_mod)
    pptx_obj = ppt_generator.generate_pptx(data)
    xlsx_bytes = xlsx_obj.getvalue() if hasattr(xlsx_obj, "getvalue") else xlsx_obj
    pptx_bytes = pptx_obj.getvalue() if hasattr(pptx_obj, "getvalue") else pptx_obj
    return {"data": data, "xlsx": xlsx_bytes, "pptx": pptx_bytes}


@pytest.fixture(scope="module")
def gbp_bundle():
    return _generate_bundle(UBER_GBP_BRIEF)


@pytest.fixture(scope="module")
def usd_bundle():
    # MANPOWER_BRIEF is also logistics_supply_chain -- same COMPLICATIONS
    # bullet, same Location Intelligence / City-Level Supply-Demand tables,
    # but a pure-USD plan. Used for the false-positive guard.
    return _generate_bundle(T.MANPOWER_BRIEF)


def _find_table(ws, header_label: str):
    """Return (header_row_idx, {col_label: col_idx}) for the row containing a
    cell value that CONTAINS ``header_label`` (substring, not exact) -- the
    Median Income header itself varies ("Median Income" vs "Median Income
    (USD)" depending on the active plan currency, see excel_v2's
    ``_income_header``)."""
    for row in ws.iter_rows():
        vals = [c.value for c in row]
        if any(isinstance(v, str) and header_label in v for v in vals):
            cols = {v: i + 1 for i, v in enumerate(vals) if isinstance(v, str)}
            return row[0].row, cols
    return None, {}


def _col_containing(cols: dict, substr: str) -> int:
    for label, idx in cols.items():
        if substr in label:
            return idx
    raise KeyError(f"no column header contains {substr!r}: {sorted(cols)}")


# ---------------------------------------------------------------------------
# Defect 1 -- Market Intelligence!F17:F22
# ---------------------------------------------------------------------------


def test_defect1_location_intelligence_median_income_uses_local_currency(gbp_bundle):
    wb = load_workbook(io.BytesIO(gbp_bundle["xlsx"]))
    ws = wb["Market Intelligence"]
    header_row, cols = _find_table(ws, "Median Income")
    assert header_row is not None, "Location Intelligence table not found"
    location_col = cols["Location"]
    income_col = _col_containing(cols, "Median Income")

    seen = {}
    for r in range(header_row + 1, header_row + 1 + len(_EXPECTED_MEDIAN_INCOME) + 3):
        loc = ws.cell(row=r, column=location_col).value
        if loc in _EXPECTED_MEDIAN_INCOME:
            seen[loc] = ws.cell(row=r, column=income_col).value

    assert set(seen) == set(_EXPECTED_MEDIAN_INCOME), (
        f"expected rows for {sorted(_EXPECTED_MEDIAN_INCOME)}, found {sorted(seen)}"
    )
    for loc, expected in _EXPECTED_MEDIAN_INCOME.items():
        assert seen[loc] == expected, (
            f"{loc}: Median Income cell = {seen[loc]!r}, expected {expected!r} "
            "(bare '$' on a GBP plan's own local-currency figure)"
        )


def test_defect1_usd_plan_median_income_unaffected(usd_bundle):
    """False-positive guard: a pure-USD plan's Median Income column must
    render EXACTLY as before the fix -- bare "$#,###", no "(CODE)" suffix."""
    wb = load_workbook(io.BytesIO(usd_bundle["xlsx"]))
    ws = wb["Market Intelligence"]
    header_row, cols = _find_table(ws, "Median Income")
    assert header_row is not None
    income_col = _col_containing(cols, "Median Income")
    location_col = cols["Location"]

    found_any = False
    for r in range(header_row + 1, header_row + 20):
        loc = ws.cell(row=r, column=location_col).value
        if not loc:
            break
        val = ws.cell(row=r, column=income_col).value
        if val in (None, ""):
            continue
        found_any = True
        assert isinstance(val, str) and val.startswith("$"), (
            f"{loc}: expected bare '$' prefix on a USD plan, got {val!r}"
        )
        assert "(" not in val, (
            f"{loc}: USD plan must not gain a '(CODE)' suffix, got {val!r}"
        )
    assert found_any, "no Median Income rows found on the USD fixture"


# ---------------------------------------------------------------------------
# Defect 2 -- Quality Intelligence!H8:H13
# ---------------------------------------------------------------------------


def test_defect2_city_level_salary_range_uses_local_currency(gbp_bundle):
    wb = load_workbook(io.BytesIO(gbp_bundle["xlsx"]))
    ws = wb["Quality Intelligence"]
    header_row, cols = _find_table(ws, "Salary Range")
    assert header_row is not None, "City-Level Supply-Demand Data table not found"
    market_col = cols["Market"]
    range_col = cols["Salary Range"]

    seen = {}
    for r in range(header_row + 1, header_row + 1 + len(_EXPECTED_SALARY_RANGE) + 3):
        mkt = ws.cell(row=r, column=market_col).value
        if mkt in _EXPECTED_SALARY_RANGE:
            seen[mkt] = ws.cell(row=r, column=range_col).value

    assert set(seen) == set(_EXPECTED_SALARY_RANGE), (
        f"expected rows for {sorted(_EXPECTED_SALARY_RANGE)}, found {sorted(seen)}"
    )
    for mkt, expected in _EXPECTED_SALARY_RANGE.items():
        assert seen[mkt] == expected, (
            f"{mkt}: Salary Range cell = {seen[mkt]!r}, expected {expected!r} "
            "(bare '$' range on a GBP plan's own local-currency figure)"
        )


def test_defect2_usd_plan_salary_range_unaffected(usd_bundle):
    """False-positive guard: a pure-USD plan's Salary Range column must
    render EXACTLY as before the fix -- bare "$lo - $hi", no "(CODE)"."""
    wb = load_workbook(io.BytesIO(usd_bundle["xlsx"]))
    ws = wb["Quality Intelligence"]
    header_row, cols = _find_table(ws, "Salary Range")
    assert header_row is not None
    market_col = cols["Market"]
    range_col = cols["Salary Range"]

    found_any = False
    for r in range(header_row + 1, header_row + 20):
        mkt = ws.cell(row=r, column=market_col).value
        if not mkt:
            break
        val = ws.cell(row=r, column=range_col).value
        if not isinstance(val, str) or val in ("—", ""):
            continue
        found_any = True
        assert val.count("$") == 2 and "(" not in val, (
            f"{mkt}: expected plain '$lo - $hi' on a USD plan, got {val!r}"
        )
    assert found_any, "no Salary Range rows found on the USD fixture"


# ---------------------------------------------------------------------------
# Defect 3 -- deck slide 2 KB benchmark string
# ---------------------------------------------------------------------------


def _slide2_paragraph_texts(pptx_bytes: bytes) -> list[str]:
    prs = Presentation(io.BytesIO(pptx_bytes))
    slide2 = prs.slides[1]
    texts = []
    for shp in slide2.shapes:
        if shp.has_text_frame:
            for p in shp.text_frame.paragraphs:
                texts.append("".join(r.text for r in p.runs))
    return texts


def test_defect3_slide2_kb_benchmark_marked_usd_on_gbp_plan(gbp_bundle):
    texts = _slide2_paragraph_texts(gbp_bundle["pptx"])
    matches = [t for t in texts if "CDL/last-mile" in t]
    assert matches, "expected the CDL/last-mile KB complication bullet on slide 2"
    text = matches[0]
    assert "US$25-US$50" in text, (
        f"expected the genuinely-USD benchmark marked 'US$', got {text!r}"
    )
    assert "£" not in text, (
        f"a USD-sourced benchmark must never be relabelled/converted to the "
        f"plan's own currency: {text!r}"
    )


def test_defect3_usd_plan_kb_benchmark_unaffected(usd_bundle):
    """False-positive guard: on a USD plan the SAME bullet must render
    EXACTLY as before -- bare '$25-$50', never 'US$'."""
    texts = _slide2_paragraph_texts(usd_bundle["pptx"])
    matches = [t for t in texts if "CDL/last-mile" in t]
    assert matches, "expected the CDL/last-mile KB complication bullet on slide 2"
    text = matches[0]
    assert "$25-$50" in text and "US$" not in text, (
        f"USD plan's own benchmark bullet must stay bare '$', got {text!r}"
    )


# ---------------------------------------------------------------------------
# Broader guard -- bundle_qa must report ZERO currency_symbol_mixing on a
# regenerated GBP non-US plan, and the USD plan must be untouched entirely.
# ---------------------------------------------------------------------------


def test_guard_zero_currency_symbol_mixing_on_gbp_plan(gbp_bundle):
    findings = run_bundle_qa(
        gbp_bundle["pptx"], gbp_bundle["xlsx"], gbp_bundle["data"]
    )
    mixing = [f for f in findings if f.get("code") == "currency_symbol_mixing"]
    assert mixing == [], f"currency_symbol_mixing findings remain: {mixing}"


def test_guard_zero_currency_symbol_mixing_on_usd_plan(usd_bundle):
    """The most important guard: a USD plan must be completely unaffected
    by the defect-1/2/3 fixes."""
    findings = run_bundle_qa(
        usd_bundle["pptx"], usd_bundle["xlsx"], usd_bundle["data"]
    )
    mixing = [f for f in findings if f.get("code") == "currency_symbol_mixing"]
    assert mixing == [], f"USD plan must never trigger currency_symbol_mixing: {mixing}"


# ---------------------------------------------------------------------------
# Defect 5 -- Quality Intelligence "City-Level Supply-Demand Data" market
# currency resolution silently defaulted to bare "$" for a BARE CITY NAME.
#
# Found by an outside probe generating 10 brief shapes against origin/main:
# a single-market UK plan ("London, United Kingdom") -- NOT a multi-country
# plan like defects 1-3 above -- still produced
# ``currency_symbol_mixing`` @ Quality Intelligence!H8: "$60,000 - $97,500"
# on a plan bundle_qa itself resolves as GBP.
#
# Root cause: gold_standard.enrich_city_level_data keys city_data by the
# CITY name alone (the token before the first comma -- "London", not
# "London, United Kingdom" or "United Kingdom"). excel_v2's
# `_mkt_currency_code` resolution passed that bare city name straight to
# plan_currency.currency_for_country(), which only maps COUNTRY
# names/codes/aliases -- so it silently returned None (bare "$") for any
# city whose name isn't ALSO a country-name substring (defects 1-3's
# locations were bare COUNTRY names like "UK"/"Australia", which resolve
# directly; "Mexico City" happens to contain "mexico" and worked by
# accident -- "London"/"Sydney"/"Toronto" do not).
#
# The fix: when a market's own bare label doesn't resolve, fall back to
# the plan's single, already-resolved currency (_get_active_currency())
# instead of silently defaulting to USD -- this codebase never carries a
# genuinely different currency per market (no per-market FX rate exists
# anywhere; plan_currency.py's own docstring: never invent a rate), so
# that is always at least as correct, and can never disagree with what
# every other cell on the same workbook calls "this plan's currency."
# ---------------------------------------------------------------------------

PEARSON_UK_EDUCATION_BRIEF: dict = {
    "client_name": "Pearson",
    "requester_name": "Test Requester",
    "requester_email": "test@joveo.com",
    "industry": "Education",
    "budget": "$150,000",
    "campaign_duration": "3-6 months",
    "hire_volume": "100+ hires",
    "work_environment": "remote",
    "locations": ["London, United Kingdom"],
    "roles": ["Curriculum Designer"],
    "target_roles": [{"title": "Curriculum Designer", "count": 100, "tier": "Salaried"}],
}

# A second, distinct market -- proves the fix generalises beyond "London"
# specifically (a different city, a different country/currency, a
# different industry). Coordinator's ask: "add at least one mixed-currency
# shape (non-US market, USD budget)" of my own.
AVIVA_AU_INSURANCE_BRIEF: dict = {
    "client_name": "Aviva Insurance AU",
    "requester_name": "Test Requester",
    "requester_email": "test@joveo.com",
    "industry": "Insurance",
    "budget": "$220,000",
    "campaign_duration": "3-6 months",
    "hire_volume": "60+ hires",
    "work_environment": "hybrid",
    "locations": ["Sydney, Australia"],
    "roles": ["Claims Adjuster"],
    "target_roles": [{"title": "Claims Adjuster", "count": 60, "tier": "Salaried"}],
}


@pytest.fixture(scope="module")
def pearson_uk_bundle():
    return _generate_bundle(PEARSON_UK_EDUCATION_BRIEF)


@pytest.fixture(scope="module")
def aviva_au_bundle():
    return _generate_bundle(AVIVA_AU_INSURANCE_BRIEF)


def test_defect5_single_market_city_only_label_resolves_plan_currency_gbp(
    pearson_uk_bundle,
):
    wb = load_workbook(io.BytesIO(pearson_uk_bundle["xlsx"]))
    ws = wb["Quality Intelligence"]
    header_row, cols = _find_table(ws, "Salary Range")
    assert header_row is not None
    market_col = cols["Market"]
    range_col = cols["Salary Range"]

    seen = {}
    for r in range(header_row + 1, header_row + 5):
        mkt = ws.cell(row=r, column=market_col).value
        if mkt:
            seen[mkt] = ws.cell(row=r, column=range_col).value
    assert "London" in seen, f"expected a London row, found {sorted(seen)}"
    assert seen["London"] == "£60,000 - £97,500 (GBP)", (
        f"London Salary Range = {seen['London']!r} -- bare '$' city-name "
        "currency-resolution gap not fixed"
    )


def test_defect5_second_market_sydney_au_resolves_plan_currency_aud(
    aviva_au_bundle,
):
    wb = load_workbook(io.BytesIO(aviva_au_bundle["xlsx"]))
    ws = wb["Quality Intelligence"]
    header_row, cols = _find_table(ws, "Salary Range")
    assert header_row is not None
    market_col = cols["Market"]
    range_col = cols["Salary Range"]

    seen = {}
    for r in range(header_row + 1, header_row + 5):
        mkt = ws.cell(row=r, column=market_col).value
        if mkt:
            seen[mkt] = ws.cell(row=r, column=range_col).value
    assert "Sydney" in seen, f"expected a Sydney row, found {sorted(seen)}"
    assert seen["Sydney"] == "A$60,000 - A$97,500 (AUD)", (
        f"Sydney Salary Range = {seen['Sydney']!r} -- bare '$' city-name "
        "currency-resolution gap not fixed"
    )


def test_defect5_guard_zero_currency_symbol_mixing_single_market_uk_plan(
    pearson_uk_bundle,
):
    findings = run_bundle_qa(
        pearson_uk_bundle["pptx"], pearson_uk_bundle["xlsx"], pearson_uk_bundle["data"]
    )
    mixing = [f for f in findings if f.get("code") == "currency_symbol_mixing"]
    assert mixing == [], f"currency_symbol_mixing findings remain: {mixing}"


def test_defect5_guard_zero_currency_symbol_mixing_second_market_au_plan(
    aviva_au_bundle,
):
    findings = run_bundle_qa(
        aviva_au_bundle["pptx"], aviva_au_bundle["xlsx"], aviva_au_bundle["data"]
    )
    mixing = [f for f in findings if f.get("code") == "currency_symbol_mixing"]
    assert mixing == [], f"currency_symbol_mixing findings remain: {mixing}"


if __name__ == "__main__":
    raise SystemExit(pytest.main([__file__, "-v"]))
