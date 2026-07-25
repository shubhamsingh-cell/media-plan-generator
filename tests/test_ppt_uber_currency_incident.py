"""Regression tests for the Uber GBP-plan bare-"$" incident.

INCIDENT: a GBP 2,000,000 plan for client Uber shipped with dollar signs
across the deck. The currency infrastructure (plan_currency /
_set_active_currency / _cur_symbol / _fmt_currency) already resolves GBP
correctly -- these were call sites bypassing it by calling the hardcoded-$
display_format.fmt_money() instead of the deck's own currency-aware
_fmt_currency() / _fmt_currency_whole().

Covers:
    1. Six hardcoded-$ call sites (ppt_generator.py market thesis, goal-gap,
       slide-6 takeaway, Push/Pull money list x2, Next Steps budget) now
       render the plan's own currency symbol, with the CPH-style whole-
       number ("never cents") rounding preserved.
    2. Slides 4 and 5 render the SAME channel dollar figures with the SAME
       symbol (they used to disagree: "$557.8K" on one, "£557.8K" on the
       other).
    3. _get_benchmarks()'s ad-platform-layer CPA/CPC rows are honestly
       flagged (and rendered) as USD benchmarks ("US$..."), matching the
       CPH row's existing honest-labelling convention, instead of a
       never-converted USD figure wearing a bare "£".
    4. A USD plan is completely unaffected (false-positive guard) -- this
       is the regression that matters most, since these call sites are
       exercised by every plan, not just international ones.

Runs under pytest, or standalone: ``python3 tests/test_ppt_uber_currency_incident.py``.
"""

from __future__ import annotations

import io
import os
import re
import sys

import pytest

PROJECT_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, PROJECT_ROOT)

import display_format as fmt  # noqa: E402
import ppt_generator as ppt  # noqa: E402
from pptx import Presentation  # noqa: E402


@pytest.fixture(autouse=True)
def _reset_active_currency():
    """_set_active_currency() stores the resolved plan currency in a
    thread-local (ppt._currency_tls) that persists across tests in the
    SAME process/thread. Several tests below set it to GBP directly
    (without going through generate_pptx, which would reset it itself);
    left dirty, that leaks into unrelated test files' tests that call
    _push_pull_split_line / _interpolate_next_steps directly and assume a
    USD default (e.g. tests/test_ppt_w3b_bundle_quality.py). Reset to the
    USD default before AND after every test in this file so ordering
    relative to other test files never matters.
    """
    ppt._set_active_currency({})
    yield
    ppt._set_active_currency({})


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
ROLES = [{"title": "Commercial Cab Driver", "count": 500, "tier": "Hourly"}]
CHANNEL_PCTS = {
    "programmatic_dsp": 35,
    "global_boards": 20,
    "niche_boards": 15,
    "social_media": 12,
    "regional_boards": 8,
    "employer_branding": 10,
}
# The two ad-platform figures that shipped in the real incident trace to the
# hardcoded static USD table (data_synthesizer.py _PLATFORM_BENCHMARKS, e.g.
# Google Ads cpc 2.90) via fuse_ad_platform_analysis -- reproduced here
# without importing data_synthesizer so the test stays offline/decoupled.
SYNTHESIZED_AD_PLATFORM = {
    "ad_platform_analysis": {
        "Google Ads": {"avg_cpc": 2.90, "avg_cpa": 48.96},
        "Meta (Facebook/Instagram)": {"avg_cpc": 1.86, "avg_cpa": 18.68},
    }
}


def _build_alloc(country: str) -> dict:
    import budget_engine

    return budget_engine.calculate_budget_allocation(
        total_budget=2_000_000,
        roles=ROLES,
        locations=[{"city": "Metro", "state": "", "country": country}],
        industry="hospitality_travel",
        channel_percentages=CHANNEL_PCTS,
        collar_type="blue",
        campaign_start_month=3,
    )


def _uber_plan(country: str, currency_budget_str: str, with_synth: bool = True) -> dict:
    data = {
        "client_name": "Uber",
        "industry": "hospitality_travel",
        "budget": currency_budget_str,
        "campaign_duration": "1-3 months",
        "hire_volume": "500+ hires",
        "campaign_goals": ["speed_to_hire", "cost_efficiency"],
        "locations": [country],
        "roles": ["Commercial Cab Driver"],
        "target_roles": ROLES,
        "_budget_allocation": _build_alloc(country),
    }
    if with_synth:
        data["_synthesized"] = dict(SYNTHESIZED_AD_PLATFORM)
    return data


def _gbp_uber_plan(with_synth: bool = True) -> dict:
    return _uber_plan("United Kingdom", "£2,000,000", with_synth=with_synth)


def _usd_uber_plan(with_synth: bool = True) -> dict:
    return _uber_plan("United States", "$2,000,000", with_synth=with_synth)


def _lean_gbp_plan(with_goals: bool = True) -> dict:
    """A content-LIGHT GBP plan (2 channels, no goal-gap/cited-2026 extras)
    that isolates the RESOLUTION card's goals-heading behavior (Fix 4) from
    the card's own pre-existing autofit/trim-trailing-content mechanism,
    which legitimately trims low-priority trailing paragraphs (including
    the goals section) off content-HEAVY fixtures like the full 8-channel
    Uber plan above -- unrelated to whether Fix 4 itself is correct."""
    import budget_engine

    alloc = budget_engine.calculate_budget_allocation(
        total_budget=150_000,
        roles=[{"title": "Registered Nurse", "count": 40, "tier": "mid"}],
        locations=[{"city": "London", "state": "", "country": "United Kingdom"}],
        industry="healthcare",
        channel_percentages={"Indeed": 60, "LinkedIn": 40},
        collar_type="white",
        campaign_start_month=9,
    )
    data = {
        "client_name": "Mercy Health",
        "industry": "healthcare",
        "budget": "£150,000",
        "locations": ["United Kingdom"],
        "roles": ["Registered Nurse"],
        "target_roles": [{"title": "Registered Nurse"}],
        "channel_categories": {"Indeed": True, "LinkedIn": True},
        "_budget_allocation": alloc,
    }
    if with_goals:
        data["campaign_goals"] = ["speed_to_hire", "cost_efficiency"]
    return data


def _all_slide_texts(pptx_bytes: bytes) -> list[str]:
    prs = Presentation(io.BytesIO(pptx_bytes))
    out: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                out.append(shape.text_frame.text)
    return out


def _bare_dollar_hits(texts: list[str]) -> list[str]:
    """Every '$' that is NOT part of a 'US$' marker."""
    hits = []
    for t in texts:
        stripped = t.replace("US$", "")
        if "$" in stripped:
            hits.append(t)
    return hits


# ---------------------------------------------------------------------------
# 1. Full-deck GBP plan: no bare "$" anywhere, plan symbol used consistently.
# ---------------------------------------------------------------------------
class TestGbpPlanNoHardcodedDollar:
    @pytest.fixture(scope="class")
    def texts(self):
        data = _gbp_uber_plan()
        pptx_bytes = ppt.generate_pptx(dict(data))
        return _all_slide_texts(pptx_bytes)

    def test_no_bare_dollar_sign_anywhere(self, texts):
        hits = _bare_dollar_hits(texts)
        assert not hits, (
            "GBP plan rendered a bare '$' (not part of a 'US$' marker) in: "
            f"{hits!r}"
        )

    def test_pound_symbol_actually_renders(self, texts):
        assert any("£" in t for t in texts), "Expected £ somewhere in a GBP plan deck"

    def test_slide2_market_thesis_uses_plan_symbol(self, texts):
        # ppt_generator.py market-thesis call site (was: fmt_money -> "$1,626/hire")
        assert any(
            "This plan projects 1230 hires at £1,626/hire" in t for t in texts
        ), texts

    def test_slide6_takeaway_uses_plan_symbol(self, texts):
        # ppt_generator.py budget-allocation takeaway call site (was: "$1,626")
        assert any(
            "with £1,626 average cost-per-hire" in t for t in texts
        ), texts

    def test_slide11_next_steps_uses_plan_symbol(self, texts):
        # ppt_generator.py _interpolate_next_steps budget_fmt call site
        assert any("£2M over 1-3 months" in t for t in texts), texts

    def test_slides_4_and_5_channel_money_match(self, texts):
        # slide 4 (Push/Pull money list) and slide 5 (channel-strategy
        # attribution chart) must render the SAME figure with the SAME
        # symbol -- these used to disagree ("$557.8K" vs "£557.8K").
        slide4_hit = any("Programmatic (DSP) £557.8K" in t for t in texts)
        slide5_hit = any("Programmatic DSP (£557.8K)" in t for t in texts)
        assert slide4_hit, f"slide 4 push/pull figure missing/wrong: {texts!r}"
        assert slide5_hit, f"slide 5 attribution figure missing/wrong: {texts!r}"

    def test_no_cents_artifact_reintroduced(self, texts):
        # The original fmt_money swap-in existed to avoid "$5,263.16"-style
        # cents. Confirm the GBP replacement never emits "£1,626.xx".
        assert not any(re.search(r"£1,626\.\d", t) for t in texts), texts

    # NOTE: "Client Goals:" heading presence (Fix 4) is NOT asserted against
    # this 8-channel fixture -- the RESOLUTION card's own pre-existing
    # autofit/trim-trailing-content mechanism (ppt_generator.py ~2023
    # _autofit_textframe) legitimately trims the goals section off this
    # content-heavy card when it overflows, independent of Fix 4's
    # correctness. See TestSlide2ResolutionCardGlyphConsistency below for a
    # leaner fixture that isolates the heading behavior itself.


# ---------------------------------------------------------------------------
# 2. Benchmark provenance flags (Fix 2): honest "US$" marking.
# ---------------------------------------------------------------------------
class TestBenchmarkProvenanceFlags:
    def test_ad_platform_cpc_and_cpa_flagged_as_usd_benchmark(self):
        data = _gbp_uber_plan()
        ppt._set_active_currency(data)
        bm = ppt._get_benchmarks("Hospitality & Travel", data)
        assert bm.get("cpc_is_usd_benchmark") is True, bm
        assert bm.get("cpa_is_usd_benchmark") is True, bm

    def test_own_plan_cph_still_not_flagged_usd(self):
        # Layer 0 (budget-engine CPH) is the plan's own localized figure --
        # must remain False (untouched by this fix).
        data = _gbp_uber_plan()
        ppt._set_active_currency(data)
        bm = ppt._get_benchmarks("Hospitality & Travel", data)
        assert bm.get("cph_is_usd_benchmark") is False, bm

    def test_slide5_benchmark_table_marks_all_three_rows_us_dollar(self):
        data = _gbp_uber_plan()
        ppt._set_active_currency(data)
        prs = Presentation()
        prs.slide_width = ppt.SLIDE_WIDTH
        prs.slide_height = ppt.SLIDE_HEIGHT
        ppt._build_slide_channel_strategy(prs, data)
        texts = [
            shape.text_frame.text
            for shape in prs.slides[0].shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        # All three rows (CPA / CPC / CPH) must be consistently marked --
        # not two of them wearing a false bare "£".
        assert any(t.startswith("US$19 - US$49") for t in texts), texts
        assert any(t.startswith("US$1.86 - US$2.90") for t in texts), texts
        assert any(t.startswith("US$1,500-US$3,500") for t in texts), texts
        # And no bare, unmarked "£" leaks in for these rows.
        assert not any(re.match(r"^£[\d]", t) for t in texts), texts


# ---------------------------------------------------------------------------
# 3. USD plan false-positive guard -- the regression that matters most.
# ---------------------------------------------------------------------------
class TestUsdPlanUnaffected:
    @pytest.fixture(scope="class")
    def texts(self):
        data = _usd_uber_plan()
        pptx_bytes = ppt.generate_pptx(dict(data))
        return _all_slide_texts(pptx_bytes)

    def test_no_pound_or_us_dollar_marker_leaks_in(self, texts):
        assert not any("£" in t for t in texts), texts
        assert not any("US$" in t for t in texts), texts

    def test_slide2_market_thesis_unchanged(self, texts):
        assert any(
            "This plan projects 1230 hires at $1,626/hire" in t for t in texts
        ), texts

    def test_slide6_takeaway_unchanged(self, texts):
        assert any("with $1,626 average cost-per-hire" in t for t in texts), texts

    def test_slide11_next_steps_unchanged(self, texts):
        assert any("$2M over 1-3 months" in t for t in texts), texts

    def test_slides_4_and_5_channel_money_match(self, texts):
        slide4_hit = any("Programmatic (DSP) $557.8K" in t for t in texts)
        slide5_hit = any("Programmatic DSP ($557.8K)" in t for t in texts)
        assert slide4_hit, texts
        assert slide5_hit, texts

    def test_benchmark_table_plain_dollar_no_us_prefix(self):
        data = _usd_uber_plan()
        ppt._set_active_currency(data)
        prs = Presentation()
        prs.slide_width = ppt.SLIDE_WIDTH
        prs.slide_height = ppt.SLIDE_HEIGHT
        ppt._build_slide_channel_strategy(prs, data)
        texts = [
            shape.text_frame.text
            for shape in prs.slides[0].shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        assert any(t.startswith("$19 - $49") for t in texts), texts
        assert any(t.startswith("$1.86 - $2.90") for t in texts), texts
        assert any(t.startswith("$1,500-$3,500") for t in texts), texts
        assert not any("US$" in t for t in texts), texts


# ---------------------------------------------------------------------------
# 4. Unit-level: the underlying helpers directly (decoupled from full-deck
#    layout, so these keep working even if slide geometry changes later).
# ---------------------------------------------------------------------------
class TestFmtCurrencyWholeHelper:
    def test_matches_fmt_money_rounding_for_usd(self):
        # _fmt_currency_whole must preserve fmt_money's exact "never cents"
        # rounding for the default/USD case -- byte-for-byte.
        ppt._set_active_currency({"locations": ["United States"]})
        assert ppt._fmt_currency_whole(5263.16) == fmt.fmt_money(5263.16) == "$5,263"

    def test_uses_plan_symbol_for_gbp(self):
        ppt._set_active_currency({"locations": ["United Kingdom"]})
        assert ppt._fmt_currency_whole(1626.43) == "£1,626"
        assert "$" not in ppt._fmt_currency_whole(1626.43)

    def test_negative_values_keep_sign_before_symbol(self):
        ppt._set_active_currency({"locations": ["United Kingdom"]})
        assert ppt._fmt_currency_whole(-1626.43) == "-£1,626"


class TestPushPullSplitLineHelper:
    def test_uses_active_currency_not_dollar(self):
        ppt._set_active_currency({"locations": ["United Kingdom"]})
        out = ppt._push_pull_split_line(
            [("Programmatic (DSP)", 557800.0)], 620000.0
        )
        assert "£557.8K" in out
        assert "$" not in out

    def test_usd_plan_unaffected(self):
        ppt._set_active_currency({"locations": ["United States"]})
        out = ppt._push_pull_split_line(
            [("Programmatic (DSP)", 557800.0)], 620000.0
        )
        assert "$557.8K" in out
        assert "£" not in out


class TestInterpolateNextStepsHelper:
    STEPS = [
        "Finalize weekly budget and success metrics",
        "Launch campaign within 2 business weeks of feed integration",
    ]

    def test_uses_active_currency_not_dollar(self):
        ppt._set_active_currency({"locations": ["United Kingdom"]})
        out = ppt._interpolate_next_steps(
            self.STEPS,
            {
                "budget": 2_000_000,
                "campaign_duration": "1-3 months",
                "client_name": "Uber",
            },
        )
        joined = " | ".join(out)
        assert "£2M" in joined
        assert "$" not in joined

    def test_usd_plan_unaffected(self):
        ppt._set_active_currency({"locations": ["United States"]})
        out = ppt._interpolate_next_steps(
            self.STEPS,
            {
                "budget": 2_000_000,
                "campaign_duration": "1-3 months",
                "client_name": "Uber",
            },
        )
        joined = " | ".join(out)
        assert "$2M" in joined
        assert "£" not in joined


# ---------------------------------------------------------------------------
# 5. Em-dash copy defects (Fix 3): raw ASCII "--" replaced with a real em
#    dash in the competitor-prose template banks.
# ---------------------------------------------------------------------------
class TestEmDashCopyDefects:
    def test_ppt_competitor_why_templates_no_ascii_double_hyphen(self):
        for tmpl in (
            ppt._COMPETITOR_WHY_TEMPLATES_INDUSTRY
            + ppt._COMPETITOR_WHY_TEMPLATES_TALENT_MARKET
        ):
            assert "--" not in tmpl, tmpl

    def test_ppt_competitor_why_templates_still_have_em_dash(self):
        # Sanity: confirm this replaced the dash glyph, not deleted it
        # (brief requires "only change the dash glyph, don't reword").
        assert any(
            "—" in t
            for t in (
                ppt._COMPETITOR_WHY_TEMPLATES_INDUSTRY
                + ppt._COMPETITOR_WHY_TEMPLATES_TALENT_MARKET
            )
        )

    def test_insight_composer_skeleton_banks_no_ascii_double_hyphen(self):
        import insight_composer

        for bucket, bank in insight_composer._SKELETON_BANKS.items():
            for sentence in bank:
                assert "--" not in sentence, f"{bucket}: {sentence!r}"

    def test_insight_composer_intensity_escalation_no_ascii_double_hyphen(self):
        import insight_composer

        out = insight_composer.compose_counter_strategy(
            "Acme Staffing",
            {"role": "Driver", "city": "Denver", "intensity": "high", "ordinal": 0},
        )
        assert "--" not in out
        assert "—" in out
        # uber_shipped_2026_07_23 fix (unsourced-competitor-claim wave):
        # "{name} has been especially aggressive here recently" asserted
        # specific, never-observed recent behaviour by the named competitor
        # -- replaced with a claim about this lane's own intensity
        # classification instead. Pin to the new text (still proving the
        # escalation sentence renders with an em dash, not a hyphen) rather
        # than the retired wording.
        assert "flagged high-intensity" in out


# ---------------------------------------------------------------------------
# 6. Bullet-glyph consistency (Fix 4), isolated from the full GBP fixture.
# ---------------------------------------------------------------------------
class TestSlide2ResolutionCardGlyphConsistency:
    def test_goals_section_has_its_own_heading(self):
        """The channel checklist (green check, 9pt) and the goals list
        (blue dot, 8pt) must no longer share one unbroken list -- a
        sub-heading now marks the transition."""
        data = _lean_gbp_plan(with_goals=True)
        pptx_bytes = ppt.generate_pptx(dict(data))
        texts = _all_slide_texts(pptx_bytes)
        resolution_card = next(
            (t for t in texts if "Nova AI Strategy:" in t and "MARKET THESIS" in t),
            None,
        )
        assert resolution_card is not None, texts
        assert "Client Goals:" in resolution_card, resolution_card
        # Heading must appear strictly between the last channel checkmark
        # section and the first goal bullet.
        goals_idx = resolution_card.index("Client Goals:")
        speed_idx = resolution_card.find("Speed to Hire")
        assert speed_idx != -1 and speed_idx > goals_idx, resolution_card

    def test_no_goals_no_heading_no_op(self):
        """When there are no campaign goals, no dangling heading is added
        (regression guard: the fix must be conditional on `if goals:`)."""
        data = _lean_gbp_plan(with_goals=False)
        pptx_bytes = ppt.generate_pptx(dict(data))
        texts = _all_slide_texts(pptx_bytes)
        resolution_card = next(
            (t for t in texts if "Nova AI Strategy:" in t and "MARKET THESIS" in t),
            None,
        )
        assert resolution_card is not None, texts
        assert "Client Goals:" not in resolution_card


if __name__ == "__main__":
    raise SystemExit(pytest.main([__file__, "-v"]))
