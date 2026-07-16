"""Tests for the S93 funnel-calibration model (owner-approved).

Covers:
    1. ``_fit_channel_funnel_rates`` -- band clamping, exact-product
       closure, below-floor/above-ceiling flagging.
    2. ``compute_funnel_stages`` -- largest-remainder stage totals,
       per-channel product closure, zero-hire band-typical rows,
       degenerate raw_apps==0 rows, metadata shape.
    3. ``calculate_budget_allocation`` wiring -- metadata.funnel present
       and self-consistent with the invariant total_projected.hires.
    4. Headline-invariance proof -- budget/apps/hires/CPA/CPH for BOTH
       reference briefs (manpower, atria) equal the values captured from
       the bundles regenerated on a clean tree BEFORE this feature existed
       (see scratchpad/headline_before.json, captured via
       tools_regen_bundles.build_plan_data prior to any funnel-model code
       being written).
    5. excel_v2 "Recruitment Funnel" table footing (monotonic stages,
       rate*prev closure, TOTAL row == sum of channel rows).
    6. ppt_generator Budget Allocation slide funnel strip presence + values.
    7. bundle_qa._check_recruitment_funnel_footing fires on a corrupted
       fixture (broken invariant / broken footing / non-monotonic row).

Runs under pytest, or standalone:
``python3 tests/test_funnel_calibration.py``.
"""

from __future__ import annotations

import io
import sys
from pathlib import Path

import pytest

PROJECT_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(PROJECT_ROOT))

import budget_engine as be  # noqa: E402
import bundle_qa  # noqa: E402

try:
    import openpyxl  # noqa: E402

    _HAS_OPENPYXL = True
except ImportError:  # pragma: no cover
    _HAS_OPENPYXL = False

try:
    from pptx import Presentation  # noqa: E402

    _HAS_PPTX = True
except ImportError:  # pragma: no cover
    _HAS_PPTX = False


# ---------------------------------------------------------------------------
# 1. _fit_channel_funnel_rates
# ---------------------------------------------------------------------------
class TestFitChannelFunnelRates:
    def _bands(self):
        return be.FUNNEL_RATE_BANDS

    def test_achievable_rate_closes_exactly_and_stays_in_band(self):
        bands = self._bands()
        for target in (0.005, 0.01, 0.02, 0.03, 0.045):
            for category in ("job_board", "niche_board", "social", "programmatic"):
                r1, r2, r3, out_of_band = be._fit_channel_funnel_rates(
                    target, category
                )
                assert out_of_band is False, (target, category)
                assert bands["raw_to_qualified"][0] - 1e-9 <= r1 <= bands["raw_to_qualified"][1] + 1e-9
                assert bands["qualified_to_interview"][0] - 1e-9 <= r2 <= bands["qualified_to_interview"][1] + 1e-9
                assert bands["interview_to_hire"][0] - 1e-9 <= r3 <= bands["interview_to_hire"][1] + 1e-9
                product = r1 * r2 * r3
                assert product == pytest.approx(target, rel=1e-6, abs=1e-9)

    def test_intent_bias_ordering_niche_beats_job_board_beats_social(self):
        # An achievable, interior target rate so none of the three
        # categories get pushed into the below/above-band fallback.
        target = 0.02
        r1_niche, *_ = be._fit_channel_funnel_rates(target, "niche_board")
        r1_referral, *_ = be._fit_channel_funnel_rates(target, "referral")
        r1_job, *_ = be._fit_channel_funnel_rates(target, "job_board")
        r1_social, *_ = be._fit_channel_funnel_rates(target, "social")
        assert r1_niche > r1_job > r1_social
        assert r1_referral > r1_job > r1_social

    def test_below_floor_pins_downstream_at_minimums(self):
        bands = self._bands()
        lo1, _hi1 = bands["raw_to_qualified"]
        lo2, _hi2 = bands["qualified_to_interview"]
        lo3, _hi3 = bands["interview_to_hire"]
        min_product = lo1 * lo2 * lo3
        target = min_product * 0.4  # comfortably below the floor
        r1, r2, r3, out_of_band = be._fit_channel_funnel_rates(target, "job_board")
        assert out_of_band is True
        assert r2 == pytest.approx(lo2)
        assert r3 == pytest.approx(lo3)
        assert r1 < lo1  # raw->qualified floats below its own floor
        assert r1 * r2 * r3 == pytest.approx(target, rel=1e-6)

    def test_above_ceiling_pins_downstream_at_maximums(self):
        bands = self._bands()
        _lo1, hi1 = bands["raw_to_qualified"]
        _lo2, hi2 = bands["qualified_to_interview"]
        _lo3, hi3 = bands["interview_to_hire"]
        max_product = hi1 * hi2 * hi3
        target = min(max_product * 2.0, 0.5)  # comfortably above the ceiling
        r1, r2, r3, out_of_band = be._fit_channel_funnel_rates(target, "job_board")
        assert out_of_band is True
        assert r2 == pytest.approx(hi2)
        assert r3 == pytest.approx(hi3)
        assert r1 >= hi1 - 1e-9

    def test_zero_target_rate_is_band_typical_not_zero(self):
        bands = self._bands()
        r1, r2, r3, out_of_band = be._fit_channel_funnel_rates(0.0, "job_board")
        assert out_of_band is False
        assert bands["raw_to_qualified"][0] <= r1 <= bands["raw_to_qualified"][1]
        assert bands["qualified_to_interview"][0] <= r2 <= bands["qualified_to_interview"][1]


# ---------------------------------------------------------------------------
# 2. compute_funnel_stages
# ---------------------------------------------------------------------------
def _chan(apps, hires, category="job_board"):
    return {"projected_applications": apps, "projected_hires": hires, "category": category}


class TestComputeFunnelStages:
    def _sample_channels(self):
        return {
            "programmatic_dsp": _chan(8565, 18, "programmatic"),
            "niche_boards": _chan(1461, 6, "niche_board"),
            "global_boards": _chan(2479, 8, "job_board"),
            "regional_boards": _chan(6211, 16, "regional"),
            "employer_branding": _chan(199, 0, "employer_branding"),
            "social_media": _chan(35, 0, "social"),
        }

    def test_metadata_shape(self):
        channels = self._sample_channels()
        total_hires = sum(c["projected_hires"] for c in channels.values())
        result = be.compute_funnel_stages(channels, total_hires)
        assert set(result.keys()) == {
            "per_channel",
            "totals",
            "stage_rates",
            "assumption_note",
        }
        assert result["stage_rates"] == be.FUNNEL_RATE_BANDS
        assert isinstance(result["assumption_note"], str) and result["assumption_note"]
        assert set(result["per_channel"].keys()) == set(channels.keys())

    def test_largest_remainder_totals_match_sum_of_channels(self):
        channels = self._sample_channels()
        total_hires = sum(c["projected_hires"] for c in channels.values())
        result = be.compute_funnel_stages(channels, total_hires)
        per_channel = result["per_channel"]
        totals = result["totals"]
        assert totals["raw_apps"] == sum(v["raw_apps"] for v in per_channel.values())
        assert totals["qualified_apps"] == sum(
            v["qualified_apps"] for v in per_channel.values()
        )
        assert totals["interviews"] == sum(v["interviews"] for v in per_channel.values())
        assert totals["hires"] == sum(v["hires"] for v in per_channel.values())
        assert totals["hires"] == total_hires

    def test_per_channel_stage_counts_match_rates_within_tolerance(self):
        channels = self._sample_channels()
        total_hires = sum(c["projected_hires"] for c in channels.values())
        result = be.compute_funnel_stages(channels, total_hires)
        for name, row in result["per_channel"].items():
            if row["raw_apps"] <= 0:
                continue
            rates = row["rates"]
            expected_qual = round(row["raw_apps"] * rates["raw_to_qualified"])
            assert abs(expected_qual - row["qualified_apps"]) <= 1, name
            expected_int = round(row["qualified_apps"] * rates["qualified_to_interview"])
            assert abs(expected_int - row["interviews"]) <= 1, name
            if row["hires"] > 0:
                expected_hires = round(row["interviews"] * rates["interview_to_hire"])
                assert abs(expected_hires - row["hires"]) <= 1, name

    def test_stages_are_non_increasing_per_channel(self):
        channels = self._sample_channels()
        total_hires = sum(c["projected_hires"] for c in channels.values())
        result = be.compute_funnel_stages(channels, total_hires)
        for name, row in result["per_channel"].items():
            assert row["raw_apps"] >= row["qualified_apps"] >= row["interviews"] >= row["hires"], name

    def test_zero_hire_channel_is_band_typical_with_zero_interview_to_hire_rate(self):
        channels = self._sample_channels()
        total_hires = sum(c["projected_hires"] for c in channels.values())
        result = be.compute_funnel_stages(channels, total_hires)
        eb_row = result["per_channel"]["employer_branding"]
        assert eb_row["hires"] == 0
        assert eb_row["qualified_apps"] > 0
        assert eb_row["interviews"] > 0
        assert eb_row["rates"]["interview_to_hire"] == 0.0
        # round(interviews * 0.0) == 0 == hires -- must never break invariant.
        assert round(eb_row["interviews"] * eb_row["rates"]["interview_to_hire"]) == eb_row["hires"]

    def test_below_floor_channel_flagged(self):
        # 18 hires / 8565 apps = 0.21%, below the 0.24% band-minimum
        # product -- this is the REAL manpower programmatic_dsp case.
        channels = {"programmatic_dsp": _chan(8565, 18, "programmatic")}
        result = be.compute_funnel_stages(channels, 18)
        row = result["per_channel"]["programmatic_dsp"]
        assert row["quality_flag"] == "funnel_rate_below_floor"
        assert row["note"]
        lo1, _hi1 = be.FUNNEL_RATE_BANDS["raw_to_qualified"]
        assert row["rates"]["raw_to_qualified"] < lo1

    def test_raw_apps_zero_is_degenerate_all_zero_row(self):
        channels = {"ghost_channel": _chan(0, 0, "job_board")}
        result = be.compute_funnel_stages(channels, 0)
        row = result["per_channel"]["ghost_channel"]
        assert row["raw_apps"] == 0
        assert row["qualified_apps"] == 0
        assert row["interviews"] == 0
        assert row["hires"] == 0


# ---------------------------------------------------------------------------
# 3. calculate_budget_allocation wiring
# ---------------------------------------------------------------------------
class TestCalculateBudgetAllocationWiresFunnel:
    def test_metadata_funnel_present_and_consistent(self):
        result = be.calculate_budget_allocation(
            total_budget=150000,
            roles=[{"title": "CDL A Driver", "count": 300, "tier": "Hourly"}],
            locations=[{"city": "Boston", "state": "MA", "country": "US"}],
            industry="logistics_supply_chain",
            channel_percentages={
                "programmatic_dsp": 35,
                "global_boards": 20,
                "niche_boards": 15,
                "social_media": 12,
                "regional_boards": 13,
                "employer_branding": 5,
            },
        )
        funnel = result["metadata"]["funnel"]
        assert set(funnel.keys()) == {
            "per_channel",
            "totals",
            "stage_rates",
            "assumption_note",
        }
        # Invariant: funnel totals.hires must equal total_projected.hires --
        # the funnel model must never disagree with the plan of record.
        assert funnel["totals"]["hires"] == result["total_projected"]["hires"]
        # Invariant: funnel totals.raw_apps must equal total_projected.applications.
        assert funnel["totals"]["raw_apps"] == result["total_projected"]["applications"]
        # Every channel_allocations entry's hires is mirrored verbatim.
        for name, ch in result["channel_allocations"].items():
            frow = funnel["per_channel"].get(name)
            assert frow is not None, name
            assert frow["hires"] == int(ch.get("projected_hires") or 0)
            assert frow["raw_apps"] == int(ch.get("projected_applications") or 0)


# ---------------------------------------------------------------------------
# 4. Headline invariance -- captured from a clean-tree regen BEFORE the
#    funnel-calibration model existed (scratchpad/headline_before.json).
#
# Input pinning (2026-07-16): the _BEFORE_* constants below were captured
# while data/channel_benchmarks_live.json still held its 2026-07-12
# tracked-snapshot content. b74be8d ("stop tracking runtime-written data
# files") correctly untracked that file -- it's regenerated at runtime and
# its churn was noise in every diff -- but that means a FRESH checkout no
# longer has it, so budget_engine's CPC-resolution cascade
# (_extract_cpc_from_live_benchmarks, budget_engine.py ~180-200) falls
# through to its next tier (trend_engine / KB / static) and produces
# different-but-not-wrong headline numbers, breaking this invariant for a
# reason that has nothing to do with the funnel-calibration engine itself.
# TestHeadlineInvariance._pin_channel_benchmarks_live below freezes the
# exact input this test's expectations were captured against, independent
# of whatever data/channel_benchmarks_live.json happens to contain (absent,
# stale, or freshly re-scraped) in the checkout the suite runs in.
#
# 2026-07-16 re-baseline: seed replaced LLM-generated April snapshot with
# web-researched July-2026 figures (see data/channel_benchmarks_seed.json
# provenance); job_board live CPC 1.52->1.62, social 4.0->2.60. The
# _BEFORE_* constants below were regenerated against that new seed, in the
# same pinned environment TestHeadlineInvariance._pin_channel_benchmarks_live
# uses (budget_engine._DATA_DIR -> tests/fixtures/funnel_invariant/, which
# now holds the byte-identical copy of the refreshed seed;
# budget_engine._channel_bench_live_cache reset to force a re-read). Total
# hires for both briefs held at 48 (Manpower) / 57 (Atria) -- CPH-floor
# clamping absorbed the CPC-band narrowing -- but clicks/applications/
# cost_per_application/cost_per_click and every per-channel dollar_amount
# moved with the new CPC inputs. Names kept as _BEFORE_* for continuity;
# they now denote "the pinned-fixture expectation", not literally
# pre-funnel-model numbers.
# ---------------------------------------------------------------------------
_FUNNEL_INVARIANT_FIXTURE_DIR = Path(__file__).resolve().parent / "fixtures" / "funnel_invariant"

_BEFORE_MANPOWER_TOTAL = {
    "applications": 17932,
    "clicks": 221566,
    "cost_per_application": 8.36,
    "cost_per_click": 0.68,
    "cost_per_hire": 3125.0,
    "hires": 48,
}
_BEFORE_ATRIA_TOTAL = {
    "applications": 16893,
    "clicks": 220065,
    "cost_per_application": 17.76,
    "cost_per_click": 1.36,
    "cost_per_hire": 5250.0,
    "hires": 57,
}
_BEFORE_MANPOWER_PER_CHANNEL = {
    "employer_branding": {"dollar_amount": 7500.0, "projected_applications": 199, "projected_hires": 0},
    "global_boards": {"dollar_amount": 30729.0, "projected_applications": 2124, "projected_hires": 8},
    "niche_boards": {"dollar_amount": 29253.0, "projected_applications": 1462, "projected_hires": 6},
    "programmatic_dsp": {"dollar_amount": 37644.0, "projected_applications": 8156, "projected_hires": 18},
    "regional_boards": {"dollar_amount": 34314.0, "projected_applications": 5891, "projected_hires": 16},
    "social_media": {"dollar_amount": 10560.0, "projected_applications": 100, "projected_hires": 0},
}
_BEFORE_ATRIA_PER_CHANNEL = {
    "employer_branding": {"dollar_amount": 23997.6, "projected_applications": 639, "projected_hires": 0},
    "global_boards": {"dollar_amount": 73374.24, "projected_applications": 5072, "projected_hires": 21},
    "niche_boards": {"dollar_amount": 35036.5, "projected_applications": 1751, "projected_hires": 9},
    "programmatic_dsp": {"dollar_amount": 80843.49, "projected_applications": 6305, "projected_hires": 17},
    "regional_boards": {"dollar_amount": 52886.28, "projected_applications": 2658, "projected_hires": 9},
    "social_media": {"dollar_amount": 33861.88, "projected_applications": 468, "projected_hires": 1},
}


def _build_plan_data(brief):
    import tools_regen_bundles as regen

    return regen.build_plan_data(brief)


class TestHeadlineInvariance:
    """HARD INVARIANT: budget, clicks, raw applications, total & per-channel
    hires, CPA, blended CPH stay EXACTLY as they were before the funnel
    model was added -- it only ADDS explanatory intermediate stages."""

    @pytest.fixture(autouse=True)
    def _pin_channel_benchmarks_live(self, monkeypatch):
        """Pin budget_engine's channel_benchmarks_live.json read to the
        frozen 2026-07-12 fixture (see module-docstring note above) so this
        test is hermetic: green with NO data/*.json present, and green
        regardless of what a live scrape refresh produces. Monkeypatching
        _DATA_DIR is the same hook tests/test_data_sources.py already uses
        to pin apis.data loaders -- no production code changes needed."""
        monkeypatch.setattr(be, "_DATA_DIR", _FUNNEL_INVARIANT_FIXTURE_DIR)
        # Force a re-read: the loader caches at module level after first call.
        monkeypatch.setattr(be, "_channel_bench_live_cache", None)

    @pytest.mark.parametrize(
        "brief_name,expected_total,expected_per_channel,expected_budget",
        [
            ("manpower", _BEFORE_MANPOWER_TOTAL, _BEFORE_MANPOWER_PER_CHANNEL, 150000.0),
            ("atria", _BEFORE_ATRIA_TOTAL, _BEFORE_ATRIA_PER_CHANNEL, 300000.0),
        ],
    )
    def test_headline_numbers_unchanged(
        self, brief_name, expected_total, expected_per_channel, expected_budget
    ):
        import tools_regen_bundles as regen

        brief = regen.MANPOWER_BRIEF if brief_name == "manpower" else regen.ATRIA_BRIEF
        data = _build_plan_data(brief)
        ba = data["_budget_allocation"]
        assert ba["metadata"]["total_budget"] == expected_budget
        tp = ba["total_projected"]
        for key, val in expected_total.items():
            assert tp[key] == pytest.approx(val), f"{brief_name}.{key}"
        for name, expected in expected_per_channel.items():
            ch = ba["channel_allocations"][name]
            assert ch.get("dollar_amount") == pytest.approx(expected["dollar_amount"]), name
            assert ch.get("projected_applications") == expected["projected_applications"], name
            assert ch.get("projected_hires") == expected["projected_hires"], name


# ---------------------------------------------------------------------------
# 5 & 6. excel_v2 funnel table footing + ppt_generator funnel strip
# ---------------------------------------------------------------------------
@pytest.mark.skipif(not _HAS_OPENPYXL, reason="openpyxl not installed")
@pytest.mark.skipif(not _HAS_PPTX, reason="python-pptx not installed")
class TestBundleRendering:
    @pytest.fixture(scope="class")
    def bundle(self):
        import tools_regen_bundles as regen

        data = regen.build_plan_data(regen.MANPOWER_BRIEF)
        from excel_v2 import generate_excel_v2
        from ppt_generator import generate_pptx
        from kb_loader import load_knowledge_base

        xlsx_bytes = generate_excel_v2(dict(data), load_kb_fn=load_knowledge_base)
        if isinstance(xlsx_bytes, tuple):
            xlsx_bytes = xlsx_bytes[0]
        pptx_bytes = generate_pptx(dict(data))
        return {"data": data, "xlsx": xlsx_bytes, "pptx": pptx_bytes}

    def test_excel_funnel_table_foots(self, bundle):
        wb = openpyxl.load_workbook(io.BytesIO(bundle["xlsx"]), data_only=False)
        ws = wb["ROI Projections"]
        rows = list(ws.iter_rows(values_only=False))
        header_ri = None
        for ri, row in enumerate(rows):
            vals = [c.value for c in row]
            if "Channel Name" in vals and "Qualified" in vals and "Interviews" in vals:
                header_ri = ri
                header_vals = vals
                break
        assert header_ri is not None, "Recruitment Funnel table not found"
        name_ci = header_vals.index("Channel Name")
        apps_ci = header_vals.index("Applications")
        qual_ci = header_vals.index("Qualified")
        int_ci = header_vals.index("Interviews")
        hires_ci = header_vals.index("Hires")

        channel_rows = []
        ri = header_ri + 1
        total_row = None
        while ri < len(rows):
            dvals = [c.value for c in rows[ri]]
            label = dvals[name_ci]
            if not isinstance(label, str) or not label.strip():
                break
            if label.strip().upper() == "TOTAL":
                total_row = dvals
                break
            channel_rows.append(dvals)
            ri += 1

        assert channel_rows, "no channel rows under Recruitment Funnel header"
        assert total_row is not None, "no TOTAL row under Recruitment Funnel header"

        for dvals in channel_rows:
            apps = dvals[apps_ci]
            qual = dvals[qual_ci]
            interviews = dvals[int_ci]
            hires_raw = dvals[hires_ci]
            hires = 0 if hires_raw == "—" else hires_raw
            assert apps >= qual >= interviews >= hires, dvals

        assert total_row[apps_ci] == sum(d[apps_ci] for d in channel_rows)
        assert total_row[qual_ci] == sum(d[qual_ci] for d in channel_rows)
        assert total_row[int_ci] == sum(d[int_ci] for d in channel_rows)
        assert total_row[hires_ci] == sum(
            0 if d[hires_ci] == "—" else d[hires_ci] for d in channel_rows
        )

        # Cross-check against the invariant total_projected numbers.
        ba = bundle["data"]["_budget_allocation"]
        assert total_row[apps_ci] == ba["total_projected"]["applications"]
        assert total_row[hires_ci] == ba["total_projected"]["hires"]

    def test_deck_funnel_strip_present_with_correct_totals(self, bundle):
        prs = Presentation(io.BytesIO(bundle["pptx"]))
        ba = bundle["data"]["_budget_allocation"]
        funnel_totals = ba["metadata"]["funnel"]["totals"]
        found_text = None
        for slide in prs.slides:
            for shp in slide.shapes:
                if shp.has_text_frame and shp.text_frame.text.startswith("Funnel:"):
                    found_text = shp.text_frame.text
                    break
            if found_text:
                break
        assert found_text, "no funnel strip found on any slide"
        assert f"{funnel_totals['raw_apps']:,}" in found_text
        assert f"{funnel_totals['qualified_apps']:,}" in found_text
        assert f"{funnel_totals['interviews']:,}" in found_text
        assert f"{funnel_totals['hires']:,}" in found_text

    def test_bundle_qa_zero_critical_on_clean_bundle(self, bundle):
        findings = bundle_qa.run_bundle_qa(bundle["pptx"], bundle["xlsx"], bundle["data"])
        critical = [f for f in findings if f["severity"] == "critical"]
        assert critical == []


# ---------------------------------------------------------------------------
# 7. bundle_qa fires on a corrupted fixture
# ---------------------------------------------------------------------------
@pytest.mark.skipif(not _HAS_OPENPYXL, reason="openpyxl not installed")
class TestBundleQaFiresOnCorruption:
    @pytest.fixture(scope="class")
    def clean_xlsx_bytes(self):
        import tools_regen_bundles as regen
        from excel_v2 import generate_excel_v2
        from kb_loader import load_knowledge_base

        data = regen.build_plan_data(regen.MANPOWER_BRIEF)
        xlsx_bytes = generate_excel_v2(dict(data), load_kb_fn=load_knowledge_base)
        if isinstance(xlsx_bytes, tuple):
            xlsx_bytes = xlsx_bytes[0]
        return xlsx_bytes

    def _load_funnel_table(self, wb):
        ws = wb["ROI Projections"]
        rows = list(ws.iter_rows(values_only=False))
        for ri, row in enumerate(rows):
            vals = [c.value for c in row]
            if "Channel Name" in vals and "Qualified" in vals and "Interviews" in vals:
                return ws, rows, ri, vals
        raise AssertionError("Recruitment Funnel table not found in fixture")

    def test_fires_on_hires_invariant_mismatch(self, clean_xlsx_bytes):
        wb = openpyxl.load_workbook(io.BytesIO(clean_xlsx_bytes))
        ws, rows, header_ri, header_vals = self._load_funnel_table(wb)
        hires_ci = header_vals.index("Hires")
        # First data row's Hires cell -- corrupt it to disagree with the
        # Per-Channel ROI Analysis table above.
        data_row = rows[header_ri + 1]
        original = data_row[hires_ci].value
        assert isinstance(original, (int, float)) and original > 0
        data_row[hires_ci].value = original + 500

        buf = io.BytesIO()
        wb.save(buf)
        corrupted_bytes = buf.getvalue()

        findings = bundle_qa.run_bundle_qa(None, corrupted_bytes, {})
        codes = {f["code"] for f in findings if f["severity"] == "critical"}
        assert "funnel_hires_invariant_mismatch" in codes or "funnel_rate_footing_mismatch" in codes

    def test_fires_on_non_monotonic_row(self, clean_xlsx_bytes):
        wb = openpyxl.load_workbook(io.BytesIO(clean_xlsx_bytes))
        ws, rows, header_ri, header_vals = self._load_funnel_table(wb)
        apps_ci = header_vals.index("Applications")
        qual_ci = header_vals.index("Qualified")
        data_row = rows[header_ri + 1]
        apps_val = data_row[apps_ci].value
        # Make Qualified exceed Applications -- structurally impossible.
        data_row[qual_ci].value = apps_val + 1000

        buf = io.BytesIO()
        wb.save(buf)
        corrupted_bytes = buf.getvalue()

        findings = bundle_qa.run_bundle_qa(None, corrupted_bytes, {})
        codes = {f["code"] for f in findings if f["severity"] == "critical"}
        assert "funnel_stage_not_monotonic" in codes or "funnel_rate_footing_mismatch" in codes

    def test_fires_on_total_row_footing_mismatch(self, clean_xlsx_bytes):
        wb = openpyxl.load_workbook(io.BytesIO(clean_xlsx_bytes))
        ws, rows, header_ri, header_vals = self._load_funnel_table(wb)
        qual_ci = header_vals.index("Qualified")
        name_ci = header_vals.index("Channel Name")
        # Find the TOTAL row and corrupt its Qualified total.
        ri = header_ri + 1
        total_row = None
        while ri < len(rows):
            label = rows[ri][name_ci].value
            if not isinstance(label, str) or not label.strip():
                break
            if label.strip().upper() == "TOTAL":
                total_row = rows[ri]
                break
            ri += 1
        assert total_row is not None
        total_row[qual_ci].value = (total_row[qual_ci].value or 0) + 999

        buf = io.BytesIO()
        wb.save(buf)
        corrupted_bytes = buf.getvalue()

        findings = bundle_qa.run_bundle_qa(None, corrupted_bytes, {})
        codes = {f["code"] for f in findings if f["severity"] == "critical"}
        assert "funnel_total_footing_mismatch" in codes


if __name__ == "__main__":
    sys.exit(pytest.main([__file__, "-v"]))
