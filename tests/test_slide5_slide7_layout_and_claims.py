"""Regression tests for the slide 5 / slide 7 text-collision defects and the
unsourced-competitor-claim brand-liability fix (real Uber bundle, GBP 2M,
"commercial cab driver", 6 non-US markets, industry misclassified as
Hospitality & Travel -- see fix/bundle-quality and UBER_BUNDLE_AUDIT.md).

Part 1 -- text collisions (all render-verified against a real Keynote
export before writing these):

  1a. Slide 5 ("Channel Strategy & Investment") benchmark table: the
      Industry CPC value ("US$0.25-US$1.00  Volatile (+53.2% MoM in July
      2025 per Recruitics)", from data/recruitment_benchmarks_deep.json's
      hospitality_travel.cpc.trend_yoy) wraps to multiple lines inside a
      row whose height/textbox were fixed regardless of content, and the
      trend text -- which already carries its own parens -- was wrapped in
      another set, producing doubled closing parens.

  1b. Slide 5 category-attribution card: 5 channels (Global Job Boards,
      Niche/Industry Boards, Regional Boards, APAC Regional, EMEA Regional)
      all group under "Job Boards"; the joined list rendered past the
      card's bottom edge.

  1c. Slide 7 ("Competitive Landscape") competitor cards: "Why:" reliably
      wraps to 2 lines at this column width, but "Counter:" was hardcoded
      at a fixed offset that assumed 1 line, so Why's second line
      overprinted Counter's first.

Existing tests/test_deck_layout.py only asserts off-canvas bounds, the 8pt
floor, and font embedding -- exactly why all three shipped uncaught. These
tests instead measure independently (via ppt_generator's own
``_estimate_lines``, applied here as an external, non-circular check
against each shape's ACTUAL generated font size and declared geometry) and
assert the estimated rendered content fits inside its own box / does not
cross its container -- not merely that the box itself is on-canvas.

Part 2 -- unsourced claims about named companies: the competitor "why this
matters" templates (_COMPETITOR_WHY_TEMPLATES_INDUSTRY /
_TALENT_MARKET) used to assert observed hiring BEHAVIOUR about a real
named company ("actively recruits", "drawing from the same", "puts direct
pressure on", "visibly hiring") for a competitor set that is, in the
common case, an inferred industry-keyed fallback with no relation to the
client's actual roles. Fixed to presence/capability framing only, plus a
provenance/hedge line on the slide when the set wasn't the client's own.
Wording is pinned identical to the companion excel_v2.py/insight_composer.py
fix (tests/test_uber_competitor_claim_fix.py's ``_INFERRED_LABEL``) and to
bundle_qa.py's new ``unsourced_competitor_claim`` rule.

Runs under pytest, or standalone: ``python3 tests/test_slide5_slide7_layout_and_claims.py``.
"""

from __future__ import annotations

import io
import sys
from pathlib import Path

PROJECT_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(PROJECT_ROOT))

import ppt_generator as ppt  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE  # noqa: E402

EMU_PER_IN = 914400
_TOL_IN = 0.02

# Same wording bundle_qa.py / excel_v2.py / insight_composer.py's companion
# fix uses (tests/test_uber_competitor_claim_fix.py's _INFERRED_LABEL) --
# pinned here too so a future edit can't silently drift the two surfaces
# out of sync.
_INFERRED_LABEL = (
    "Competitor set inferred from industry classification; not verified "
    "against live posting data."
)

# The same asserted-behaviour verb phrases the companion insight_composer.py/
# bundle_qa.py fix bans (tests/test_uber_competitor_claim_fix.py's
# _BANNED_SUBSTRINGS), plus the two additional phrasings the ORIGINAL
# ppt_generator.py templates used ("visibly hiring", "pulls from the same").
_BANNED_SUBSTRINGS = (
    "actively recruit",
    "actively staffing",
    "actively competing",
    "keep pressure on",
    "keeps pressure on",
    "puts direct pressure on",
    "put direct pressure on",
    "drawing from the same",
    "draws from the same",
    "pulls from the same",
    "is slower to respond",
    "has been especially aggressive",
    "visibly hiring",
)


# ---------------------------------------------------------------------------
# Fixture: reproduces the real shipped-bundle shape
# ---------------------------------------------------------------------------
def _uber_style_plan() -> dict:
    """A rideshare client run under ``hospitality_travel`` (the wizard has
    no Rideshare/Gig card -- see fix/bundle-quality), all 8 default
    channels on for a non-US market so 5 group under "Job Boards"
    (visual:uber#2/#3), and no client-supplied competitor list so the
    industry-keyed gold-standard fallback (visual:uber#4/#5) populates
    slide 7 -- the exact real-world precondition for all four defects."""
    return {
        "client_name": "Uber",
        "industry": "hospitality_travel",
        "industry_label": "Hospitality & Travel",
        "locations": ["London, United Kingdom"],
        "roles": ["Commercial Cab Driver"],
        "budget": "£2,000,000",
        "work_environment": "hybrid",
        "channel_categories": {
            "programmatic_dsp": True,
            "global_boards": True,
            "niche_boards": True,
            "regional_boards": True,
            "social_media": True,
            "employer_branding": True,
            "apac_regional": True,
            "emea_regional": True,
        },
        "_gold_standard": {
            "competitors_list": ["Marriott", "Hilton", "Hyatt"],
        },
    }


def _deck() -> bytes:
    return ppt.generate_pptx(_uber_style_plan())


# ---------------------------------------------------------------------------
# Shape-geometry helpers
# ---------------------------------------------------------------------------
def _iter_shapes(shapes):
    for sh in shapes:
        yield sh
        if sh.shape_type == 6:  # group
            try:
                yield from _iter_shapes(sh.shapes)
            except Exception:
                pass


def _slide_by_headline(prs: Presentation, headline: str):
    """Locate a slide by its big top-band title (26pt, top=0.45in) -- NOT
    a plain substring search, since e.g. "Competitive Landscape" also
    appears as a Risk Analysis card title on a different slide."""
    for slide in prs.slides:
        for sh in _iter_shapes(slide.shapes):
            if not (sh.has_text_frame and sh.text_frame.text.strip()):
                continue
            if sh.text_frame.text.strip() != headline:
                continue
            top_in = sh.top / EMU_PER_IN
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.size and abs(r.font.size.pt - 26.0) < 0.1 and 0.4 < top_in < 0.5:
                        return slide
    raise AssertionError(f"no slide headline found matching {headline!r}")


def _text_shapes(slide):
    return [
        sh
        for sh in _iter_shapes(slide.shapes)
        if sh.has_text_frame and sh.text_frame.text.strip()
    ]


def _run_font_pt(shape, default: float = 10.0) -> float:
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if r.font.size is not None:
                return r.font.size.pt
    return default


def _content_h_in(text: str, width_in: float, font_pt: float, char_em: float = 0.53) -> float:
    """Independent (test-owned) estimate of rendered content height, using
    ppt_generator's own module-level ``_estimate_lines`` (the established
    line-wrap estimator, reused rather than re-implemented) against the
    line-height factor _autofit_textframe documents as calibrated to the
    real Keynote render ("~1.4x the point size once leading ... is
    included"). This is a real, non-circular check: it reads the shape's
    ACTUAL generated font size/width/height from the produced .pptx, not a
    hardcoded expectation, and is a DIFFERENT computation path (own
    line-height application, and -- where passed -- its own char-width
    constant) than whatever the fix used internally to decide the box's
    own height. Default ``char_em`` matches ``_estimate_lines``'s own
    default (average body prose); callers pass a wider value for
    demonstrably wider content classes (see the CPC-row test below).
    """
    n_lines = ppt._estimate_lines(text, width_in, font_pt, char_em=char_em)
    return n_lines * (font_pt * 1.42) / 72.0


def _rounded_rect_cards(slide, min_w_in: float, min_h_in: float):
    """Card-sized rounded-rectangle background shapes (no text) on a slide,
    sorted by top. Distinguishes real cards from smaller rounded rects used
    elsewhere (e.g. the channel-mix bar fills)."""
    cards = []
    for sh in _iter_shapes(slide.shapes):
        if sh.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            continue
        try:
            if sh.auto_shape_type != MSO_SHAPE.ROUNDED_RECTANGLE:
                continue
        except Exception:
            continue
        if sh.has_text_frame and sh.text_frame.text.strip():
            continue
        if sh.width / EMU_PER_IN < min_w_in or sh.height / EMU_PER_IN < min_h_in:
            continue
        cards.append(sh)
    cards.sort(key=lambda s: s.top)
    return cards


# ---------------------------------------------------------------------------
# 1a -- Slide 5 benchmark table: doubled parens + row/value containment
# ---------------------------------------------------------------------------
def test_slide5_cpc_value_has_no_doubled_parentheses():
    prs = Presentation(io.BytesIO(_deck()))
    slide = _slide_by_headline(prs, "Channel Strategy & Investment")
    cpc_cells = [
        sh.text_frame.text
        for sh in _text_shapes(slide)
        if "Volatile" in sh.text_frame.text and "Recruitics" in sh.text_frame.text
    ]
    assert cpc_cells, "expected the Industry CPC value cell (KB trend_yoy text)"
    for t in cpc_cells:
        assert "((" not in t and "))" not in t, f"doubled parentheses in {t!r}"
        # the un-doubled form still carries the KB's own single parens
        assert "(+53.2%" in t


def test_slide5_cpc_value_fits_inside_its_own_row():
    """visual:uber#1/#2: this value used to be laid out in a fixed
    row_h=0.38in row regardless of content; it wraps to several lines and
    overflowed into the row above/below (clipped by the next row's opaque
    background rect)."""
    prs = Presentation(io.BytesIO(_deck()))
    slide = _slide_by_headline(prs, "Channel Strategy & Investment")
    cpc_shape = next(
        (
            sh
            for sh in _text_shapes(slide)
            if "Volatile" in sh.text_frame.text and "Recruitics" in sh.text_frame.text
        ),
        None,
    )
    assert cpc_shape is not None
    width_in = cpc_shape.width / EMU_PER_IN
    height_in = cpc_shape.height / EMU_PER_IN
    font_pt = _run_font_pt(cpc_shape)
    # This value is bold currency/percent-heavy text ("US$0.25-US$1.00
    # Volatile (+53.2% ..."); a real Keynote render of it (checked while
    # writing this fix) wraps wider than the average-body-prose char_em=
    # 0.53 default predicts (0.53 estimates 2 lines at 10pt/2.5in; the
    # real render shows 3). ppt_generator.py itself uses a fatter advance
    # for this content class (_KPI_CHAR_EM=0.62, "big, bold, digit-heavy
    # ... numerals"); 0.56 is a smaller, independently-chosen bump that
    # matches the real render for this specific string.
    needed_in = _content_h_in(cpc_shape.text_frame.text, width_in, font_pt, char_em=0.56)
    assert needed_in <= height_in + _TOL_IN, (
        f"CPC value needs ~{needed_in:.2f}in of height but its own row box is "
        f"only {height_in:.2f}in tall ({cpc_shape.text_frame.text!r}) -- it "
        f"will overflow into the neighbouring row"
    )


def test_slide5_benchmark_rows_do_not_overlap_each_other():
    """Belt-and-suspenders position check on the new cascading row layout:
    each row's box must end at or before the next row's box begins."""
    prs = Presentation(io.BytesIO(_deck()))
    slide = _slide_by_headline(prs, "Channel Strategy & Investment")
    # Benchmark-table VALUE cells only (table_left=7.5in + Inches(2.4) col
    # offset =9.9in) -- restricted to one column so a metric-label cell and
    # its own row's value cell (same top, different X) don't register as a
    # false "overlap" under a Y-only comparison.
    value_cells = [
        sh
        for sh in _text_shapes(slide)
        if 9.7 < sh.left / EMU_PER_IN < 10.1 and 1.9 < sh.top / EMU_PER_IN < 5.0
    ]
    assert len(value_cells) >= 3, "expected several benchmark value cells"
    value_cells.sort(key=lambda s: s.top)
    for a, b in zip(value_cells, value_cells[1:]):
        a_bottom_in = a.top / EMU_PER_IN + a.height / EMU_PER_IN
        b_top_in = b.top / EMU_PER_IN
        assert a_bottom_in <= b_top_in + _TOL_IN, (
            f"benchmark row {a.text_frame.text!r} (bottom {a_bottom_in:.2f}in) "
            f"overlaps the next row {b.text_frame.text!r} (top {b_top_in:.2f}in)"
        )


# ---------------------------------------------------------------------------
# 1b -- Slide 5 category card: channel-list overflow past the card boundary
# ---------------------------------------------------------------------------
def test_slide5_job_boards_channel_list_stays_inside_its_card():
    """visual:uber#3: 5 channels grouped under "Job Boards" joined into one
    line used to overflow a fixed Inches(0.42) box inside a fixed
    Inches(1.2) card -- the last word rendered below the card's rounded
    bottom edge. The fix sizes the channel-list textbox's own declared
    height to its (possibly "+N more"-capped) content and caps that height
    to what's actually left inside the card -- so checking the rendered
    text's estimated content height against the box's OWN declared height
    directly verifies "does not cross its container": pre-fix the box was
    a content-blind constant; post-fix it is provably content-derived and
    card-bounded (see the enclosing-card sanity check below)."""
    prs = Presentation(io.BytesIO(_deck()))
    slide = _slide_by_headline(prs, "Channel Strategy & Investment")
    ch_list_shape = next(
        (
            sh
            for sh in _text_shapes(slide)
            if sh.text_frame.text.startswith("Global Job Boards,")
        ),
        None,
    )
    assert ch_list_shape is not None, "expected the Job Boards card's channel list"

    # Sanity: the shape really does sit inside one of the 4 category cards
    # (not floating off the attribution band entirely).
    cards = _rounded_rect_cards(slide, min_w_in=2.0, min_h_in=1.0)
    assert cards, "expected the 4 category-attribution cards"
    text_left_in = ch_list_shape.left / EMU_PER_IN
    text_top_in = ch_list_shape.top / EMU_PER_IN
    card = next(
        (
            c
            for c in cards
            if c.left / EMU_PER_IN - 0.05 <= text_left_in
            and text_left_in <= c.left / EMU_PER_IN + c.width / EMU_PER_IN
            and c.top / EMU_PER_IN - 0.05 <= text_top_in
        ),
        None,
    )
    assert card is not None, "could not find the enclosing card for the channel list"

    font_pt = _run_font_pt(ch_list_shape)
    width_in = ch_list_shape.width / EMU_PER_IN
    height_in = ch_list_shape.height / EMU_PER_IN
    needed_in = _content_h_in(ch_list_shape.text_frame.text, width_in, font_pt)
    assert needed_in <= height_in + _TOL_IN, (
        f"channel list {ch_list_shape.text_frame.text!r} needs ~{needed_in:.2f}in "
        f"of height but its own box is only {height_in:.2f}in tall -- text will "
        f"cross the card boundary"
    )


# ---------------------------------------------------------------------------
# 1c -- Slide 7 competitor cards: Why/Counter collision + card overlap
# ---------------------------------------------------------------------------
def test_slide7_why_text_fits_before_counter_starts():
    """visual:uber#4: "Why:" sat in a fixed cy+0.3in / Inches(0.25) box
    (~1 line at 8pt) while "Counter:" was hardcoded at cy+0.58in. The Why
    sentence reliably wraps to 2 lines at this column width, so its
    second line overprinted Counter's first line on every card."""
    prs = Presentation(io.BytesIO(_deck()))
    slide = _slide_by_headline(prs, "Competitive Landscape")
    why_shapes = [sh for sh in _text_shapes(slide) if sh.text_frame.text.startswith("Why:")]
    counter_shapes = [
        sh for sh in _text_shapes(slide) if sh.text_frame.text.startswith("Counter:")
    ]
    assert why_shapes and len(why_shapes) == len(counter_shapes) >= 2, (
        "expected multiple competitor cards, each with a Why: and Counter: line"
    )
    why_shapes.sort(key=lambda s: s.top)
    counter_shapes.sort(key=lambda s: s.top)
    for why_sh, counter_sh in zip(why_shapes, counter_shapes):
        width_in = why_sh.width / EMU_PER_IN
        font_pt = _run_font_pt(why_sh)
        why_needed_in = _content_h_in(why_sh.text_frame.text, width_in, font_pt)
        why_bottom_in = why_sh.top / EMU_PER_IN + why_needed_in
        counter_top_in = counter_sh.top / EMU_PER_IN
        assert why_bottom_in <= counter_top_in + _TOL_IN, (
            f"Why text {why_sh.text_frame.text!r} needs to reach "
            f"{why_bottom_in:.2f}in but Counter starts at {counter_top_in:.2f}in "
            f"-- Why's wrapped line(s) overprint Counter"
        )


def test_slide7_competitor_cards_do_not_overlap_each_other():
    prs = Presentation(io.BytesIO(_deck()))
    slide = _slide_by_headline(prs, "Competitive Landscape")
    cards = _rounded_rect_cards(slide, min_w_in=4.0, min_h_in=1.0)
    assert len(cards) >= 2, "expected multiple competitor cards"
    for a, b in zip(cards, cards[1:]):
        a_bottom_in = a.top / EMU_PER_IN + a.height / EMU_PER_IN
        b_top_in = b.top / EMU_PER_IN
        assert a_bottom_in <= b_top_in + _TOL_IN, (
            f"competitor card at top={a.top / EMU_PER_IN:.2f}in "
            f"(bottom {a_bottom_in:.2f}in) overlaps the next card "
            f"(top {b_top_in:.2f}in)"
        )


# ---------------------------------------------------------------------------
# Part 2 -- unsourced claims about named companies
# ---------------------------------------------------------------------------
def test_slide7_shows_provenance_hedge_for_inferred_competitor_set():
    """The competitor set here came from the industry-keyed gold-standard
    fallback, not the client's own brief -- the slide must disclose that."""
    prs = Presentation(io.BytesIO(_deck()))
    slide = _slide_by_headline(prs, "Competitive Landscape")
    texts = [sh.text_frame.text for sh in _text_shapes(slide)]
    assert any(_INFERRED_LABEL in t for t in texts), (
        "expected the provenance/hedge line since the competitor set is "
        "inferred, not client-supplied or independently verified"
    )


def test_slide7_no_inferred_label_when_brief_supplied_real_competitors():
    """False-positive guard: when the client's OWN brief names real
    competitors, the "inferred" disclosure must not appear."""
    data = _uber_style_plan()
    data["competitors"] = ["Bolt", "Lyft"]
    data.pop("_gold_standard", None)
    prs = Presentation(io.BytesIO(ppt.generate_pptx(data)))
    slide = _slide_by_headline(prs, "Competitive Landscape")
    texts = [sh.text_frame.text for sh in _text_shapes(slide)]
    assert not any(_INFERRED_LABEL in t for t in texts)
    assert any("Bolt" in t for t in texts)


def test_slide7_why_text_has_no_banned_asserted_behaviour_verbs():
    """copy: the Why templates must never assert OBSERVED hiring behaviour
    about a named third party -- only presence/capability framing."""
    prs = Presentation(io.BytesIO(_deck()))
    slide = _slide_by_headline(prs, "Competitive Landscape")
    why_texts = [
        sh.text_frame.text for sh in _text_shapes(slide) if sh.text_frame.text.startswith("Why:")
    ]
    assert why_texts, "expected Why: lines on the competitor cards"
    for t in why_texts:
        low = t.lower()
        for banned in _BANNED_SUBSTRINGS:
            assert banned not in low, f"banned asserted-behaviour phrase {banned!r} in {t!r}"
        # still names the competitor -- the fix must not have gone so far
        # as to stop naming it (presence/capability framing is allowed and
        # expected, not silence).
        assert any(name in t for name in ("Marriott", "Hilton", "Hyatt"))


def test_why_templates_module_level_contain_no_banned_verbs_at_any_ordinal():
    """Direct check against every template string in both buckets (not
    just what one fixture's ordinal happened to select), formatted the
    same way ``_compose_competitor_why`` does."""
    for templates in (
        ppt._COMPETITOR_WHY_TEMPLATES_INDUSTRY,
        ppt._COMPETITOR_WHY_TEMPLATES_TALENT_MARKET,
    ):
        assert len(templates) >= 4
        for tmpl in templates:
            rendered = tmpl.format(
                name="Marriott",
                pool="commercial cab driver candidates in UK",
                type_phrase="",
                article="a",
            ).lower()
            for banned in _BANNED_SUBSTRINGS:
                assert banned not in rendered, f"{banned!r} in {rendered!r}"


if __name__ == "__main__":
    import pytest as _pytest

    raise SystemExit(_pytest.main([__file__, "-v"]))


def test_slide7_competitor_card_title_has_no_internal_scope_tag():
    """The card heading must not leak the internal "(National)" scope marker.

    gold_standard.py:2927 prepends "(National) " to industry-level (non
    city-specific) competitors. ppt_generator already strips it before
    interpolating the Why/Counter prose (_strip_competitor_tag), but the
    card TITLE rendered comp_name raw -- so a real client deck showed
    "(National) Marriott" as the heading while the sentence directly
    beneath it said "Marriott". Feed tagged names the way the real
    gold-standard pipeline does, and assert the tag never reaches a slide.
    """
    plan = _uber_style_plan()
    plan["_gold_standard"] = {
        "competitor_mapping": {
            "London, United Kingdom": {
                "top_employers": [
                    "(National) Marriott",
                    "(National) Hilton",
                    "(National) Hyatt",
                ]
            }
        },
        "competitors_list": ["(National) Marriott", "(National) Hilton"],
    }
    prs = Presentation(io.BytesIO(ppt.generate_pptx(plan)))
    slide = _slide_by_headline(prs, "Competitive Landscape")
    leaked = [
        sh.text_frame.text
        for sh in _text_shapes(slide)
        if any(t in sh.text_frame.text for t in ("(National)", "(Regional)", "(Local)"))
    ]
    assert not leaked, f"internal scope tag leaked onto slide 7: {leaked}"

    # and the competitor is still actually named -- guards a fix that
    # "passes" by dropping the card entirely.
    all_text = " ".join(sh.text_frame.text for sh in _text_shapes(slide))
    assert "Marriott" in all_text, "competitor name vanished along with the tag"
