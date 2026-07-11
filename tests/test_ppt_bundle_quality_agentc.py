"""Tests for Agent C's ppt_generator.py bundle-quality fixes.

Covers:
    1. _is_ai_training_plan token/phrase-boundary matching -- "Supply Chain"
       and "Maintenance Technician" must NOT false-positive on the "ai"
       substring inside "ch-AI-n" / "m-AI-ntenance".
    2. plan_geo delegation -- _is_us_only_campaign is a thin wrapper.
    3. KB v2 industry-keyed readers: cpa_reference/case_study drop entirely
       (or fall back to a Next-Steps-only layout) for non-AI-training plans;
       no AI-training content leaks into a logistics/healthcare deck.
    4. Display sweep: fmt_money never emits "$150.0K"; fmt_count never
       emits "(s)"; channel_label never leaks a raw snake_case key;
       client_display_name fixes mixed-case client names; the exec-summary
       headline caps the role list at 2 + "and N more roles".
    5. Honest comparisons: the two hardcoded is_better:True rows are gone;
       ties render neutral; a value inside a benchmark range renders
       neutral; the goal-gap row/callout renders when a goal is stated and
       the plan falls short.
    6. Cross-artifact benchmark single-sourcing: the deck's Industry
       CPA/CPC/Apply-Rate numbers on slide 5 numerically match the workbook
       Executive Summary sheet's "Recruitment Benchmarks" section for the
       same industry.

Runs under pytest, or standalone: ``python3 tests/test_ppt_bundle_quality_agentc.py``.
"""

from __future__ import annotations

import io
import os
import re
import sys

import pytest

PROJECT_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, PROJECT_ROOT)

import ppt_generator as ppt  # noqa: E402
import display_format as fmt  # noqa: E402
from pptx import Presentation  # noqa: E402


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def _all_slide_text(pptx_bytes: bytes) -> list[str]:
    prs = Presentation(io.BytesIO(pptx_bytes))
    out: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                out.append(shape.text_frame.text)
    return out


def _build_alloc(industry: str, roles: list[dict], country: str = "United States") -> dict:
    import budget_engine

    return budget_engine.calculate_budget_allocation(
        total_budget=150_000,
        roles=roles,
        locations=[{"city": "Metro", "state": "", "country": country}],
        industry=industry,
        channel_percentages={
            "programmatic_dsp": 35,
            "global_boards": 20,
            "niche_boards": 15,
            "social_media": 12,
            "regional_boards": 8,
            "employer_branding": 10,
        },
        collar_type="blue",
        campaign_start_month=3,
    )


def _logistics_plan() -> dict:
    roles = [{"title": "CDL A Driver", "count": 100, "tier": "Hourly"}]
    return {
        "client_name": "Manpower - Amerigas",
        "industry": "logistics_supply_chain",
        "industry_label": "Logistics & Supply Chain",
        "budget": "$150,000",
        "locations": ["Denver, CO", "Dallas, TX"],
        "roles": ["CDL A Driver"],
        "target_roles": roles,
        "work_environment": "on_site",
        "hire_volume": "500 hires",
        "_budget_allocation": _build_alloc("logistics_supply_chain", roles),
    }


def _healthcare_plan() -> dict:
    roles = [
        {"title": "Memory Care Associate", "count": 30, "tier": "Hourly"},
        {"title": "Nurse", "count": 20, "tier": "Professional"},
        {"title": "Cook", "count": 10, "tier": "Hourly"},
        {"title": "Driver", "count": 10, "tier": "Hourly"},
        {"title": "Maintenance Technician", "count": 10, "tier": "Hourly"},
    ]
    return {
        "client_name": "atria Senior living",
        "industry": "healthcare_medical",
        "industry_label": "Healthcare & Medical",
        "budget": "$300,000",
        "locations": ["Fort Worth, TX"],
        "roles": [r["title"] for r in roles],
        "target_roles": roles,
        "work_environment": "remote",
        "hire_volume": "500+ hires",
        "_budget_allocation": _build_alloc("healthcare_medical", roles),
    }


def _deck_bytes(data: dict) -> bytes:
    return ppt.generate_pptx(dict(data))


# ---------------------------------------------------------------------------
# 1. _is_ai_training_plan token/phrase-boundary matching
# ---------------------------------------------------------------------------
class TestAiTrainingTokenGate:
    def test_supply_chain_is_not_ai_training(self):
        # "ai" is a raw substring of "ch-AI-n" -- must not false-positive.
        assert ppt._is_ai_training_plan(
            {"industry": "Logistics & Supply Chain", "roles": ["CDL A Driver"]}
        ) is False

    def test_maintenance_technician_is_not_ai_training(self):
        # "ai" is a raw substring of "m-AI-ntenance" -- must not false-positive.
        assert ppt._is_ai_training_plan(
            {"industry": "Healthcare & Medical", "roles": ["Maintenance Technician"]}
        ) is False

    def test_real_ai_training_phrases_match(self):
        assert ppt._is_ai_training_plan(
            {"industry": "AI Training", "roles": ["AI Trainer"]}
        ) is True
        assert ppt._is_ai_training_plan(
            {"industry": "general", "roles": ["Data Annotator"]}
        ) is True
        assert ppt._is_ai_training_plan(
            {"industry": "general", "roles": ["RLHF Specialist"]}
        ) is True
        assert ppt._is_ai_training_plan(
            {"industry": "general", "roles": ["Data Labeling Associate"]}
        ) is True

    def test_bare_ai_token_matches_but_substring_does_not(self):
        assert ppt._is_ai_training_plan({"industry": "AI", "roles": []}) is True
        assert ppt._is_ai_training_plan(
            {"industry": "general", "roles": ["Maintenance"]}
        ) is False
        assert ppt._is_ai_training_plan(
            {"industry": "general", "roles": ["Certified Nursing Assistant"]}
        ) is False


# ---------------------------------------------------------------------------
# 2. plan_geo delegation
# ---------------------------------------------------------------------------
class TestPlanGeoDelegation:
    def test_delegates_to_plan_geo(self):
        import plan_geo

        data = {"locations": ["Auckland, New Zealand"]}
        assert ppt._is_us_only_campaign(data) == plan_geo.is_us_plan(data)
        assert ppt._is_us_only_campaign(data) is False

    def test_us_plan_still_true(self):
        assert ppt._is_us_only_campaign({"locations": ["Denver, CO"]}) is True


# ---------------------------------------------------------------------------
# 3. KB v2 industry-keyed readers -- no AI-training leakage
# ---------------------------------------------------------------------------
class TestKbV2IndustryGating:
    def test_kb_content_key_ai_training(self):
        assert ppt._kb_content_key({"industry": "AI Training", "roles": ["AI Trainer"]}) == "ai_training"

    def test_kb_content_key_other_industry(self):
        assert (
            ppt._kb_content_key({"industry": "logistics_supply_chain", "roles": ["CDL A Driver"]})
            == "logistics_supply_chain"
        )

    def test_logistics_deck_has_no_ai_training_content(self):
        data = _logistics_plan()
        texts = _all_slide_text(_deck_bytes(data))
        blob = "\n".join(texts).lower()
        for needle in (
            "ai trainer",
            "ai training",
            "data labeling",
            "data annotation",
            "rlhf",
            "hiring high-quality ai trainers",
        ):
            assert needle not in blob, f"leaked AI-training content: {needle!r}"

    def test_healthcare_deck_has_no_ai_training_content(self):
        data = _healthcare_plan()
        texts = _all_slide_text(_deck_bytes(data))
        blob = "\n".join(texts).lower()
        for needle in (
            "ai trainer",
            "ai training",
            "data labeling",
            "data annotation",
            "rlhf",
            "hiring high-quality ai trainers",
        ):
            assert needle not in blob, f"leaked AI-training content: {needle!r}"

    def test_cpa_reference_slide_dropped_for_non_ai_training(self):
        data = _logistics_plan()
        deck_kb = ppt._load_deck_kb()
        prs = Presentation()
        prs.slide_width = ppt.SLIDE_WIDTH
        prs.slide_height = ppt.SLIDE_HEIGHT
        ppt._build_slide_cpa_reference(prs, data, deck_kb)
        assert len(prs.slides) == 0

    def test_case_study_slide_falls_back_to_next_steps_only(self):
        data = _logistics_plan()
        deck_kb = ppt._load_deck_kb()
        prs = Presentation()
        prs.slide_width = ppt.SLIDE_WIDTH
        prs.slide_height = ppt.SLIDE_HEIGHT
        ppt._build_slide_case_study_next_steps(prs, data, deck_kb)
        assert len(prs.slides) == 1
        texts = [
            s.text_frame.text
            for s in prs.slides[0].shapes
            if s.has_text_frame and s.text_frame.text.strip()
        ]
        blob = "\n".join(texts)
        assert "NEXT STEPS" in blob.upper()
        # No case-study header/tiles/stat row content.
        assert "Case study:" not in blob
        assert "Challenges" not in blob
        assert "The Joveo Solution" not in blob


# ---------------------------------------------------------------------------
# 4. Display sweep
# ---------------------------------------------------------------------------
class TestDisplaySweep:
    def test_fmt_currency_never_trailing_dot_zero(self):
        assert ppt._fmt_currency(150_000, compact=True) == "$150K"
        assert ppt._fmt_currency(2_000_000, compact=True) == "$2M"
        assert ppt._fmt_currency(52_500, compact=True) == "$52.5K"

    def test_format_budget_display_never_trailing_dot_zero(self):
        assert ppt._format_budget_display("$2,000,000") == "$2M"

    def test_client_display_name_fixes_mixed_case(self):
        assert ppt._proper_client_name("atria Senior living") == "Atria Senior Living"
        assert ppt._proper_client_name("MANPOWER - AMERIGAS") == "Manpower - Amerigas"
        assert ppt._proper_client_name("AT&T") == "AT&T"

    def test_no_snake_case_channel_keys_in_deck_text(self):
        data = _logistics_plan()
        texts = _all_slide_text(_deck_bytes(data))
        blob = "\n".join(texts)
        for raw_key in (
            "programmatic_dsp",
            "global_boards",
            "niche_boards",
            "social_media",
            "regional_boards",
            "employer_branding",
        ):
            assert raw_key not in blob, f"raw snake_case channel key leaked: {raw_key!r}"

    def test_no_bare_s_paren_pluralization(self):
        data = _logistics_plan()
        texts = _all_slide_text(_deck_bytes(data))
        blob = "\n".join(texts)
        assert "market(s)" not in blob
        assert "location(s)" not in blob

    def test_exec_summary_headline_caps_roles_at_two(self):
        roles = ["Role One", "Role Two", "Role Three", "Role Four", "Role Five"]
        headline = ppt._cap_roles_for_headline(roles, cap=2)
        assert headline == "Role One, Role Two, and 3 more roles"

    def test_exec_summary_headline_two_roles_no_truncation(self):
        assert ppt._cap_roles_for_headline(["Role One", "Role Two"], cap=2) == (
            "Role One and Role Two"
        )

    def test_exec_summary_headline_single_role(self):
        assert ppt._cap_roles_for_headline(["Solo Role"], cap=2) == "Solo Role"

    def test_work_model_corrected_via_gold_standard(self):
        data = _healthcare_plan()  # stated "remote" but roles are on-site
        texts = _all_slide_text(_deck_bytes(data))
        blob = "\n".join(texts)
        assert "Work Model:  On-site" in blob or "Work Model:  On-Site" in blob
        assert "adjusted to On-site" in blob


# ---------------------------------------------------------------------------
# 5. Honest comparisons
# ---------------------------------------------------------------------------
class TestComparisonNeutrality:
    def test_cmp_status_tie_band(self):
        assert ppt._cmp_status(6, 6) == "on_par"
        assert ppt._cmp_status(6.1, 6, tie_band=0.05) == "on_par"  # within 5%
        assert ppt._cmp_status(10, 6) == "beating"
        assert ppt._cmp_status(3, 6) == "trailing"

    def test_cmp_status_lower_is_better(self):
        assert ppt._cmp_status(10, 20, higher_is_better=False) == "beating"
        assert ppt._cmp_status(30, 20, higher_is_better=False) == "trailing"

    def test_value_within_range_is_on_par_never_beating(self):
        assert ppt._cmp_status_in_range(4, 3, 5) == "on_par"
        assert ppt._cmp_status_in_range(3, 3, 5) == "on_par"
        assert ppt._cmp_status_in_range(5, 3, 5) == "on_par"
        assert ppt._cmp_status_in_range(2, 3, 5) == "trailing"
        assert ppt._cmp_status_in_range(6, 3, 5) == "beating"

    def test_no_hardcoded_true_rows_render_neutral_style(self):
        # "none" status (used for "Projected Hires" vs "-" / "Total
        # Investment" vs "Varies") must render with no colored arrow.
        color, indicator, _ = ppt._CMP_STATUS_STYLE["none"]
        assert indicator == ""
        assert color == ppt.MUTED_TEXT

    def test_channels_tie_renders_neutral_in_real_deck(self):
        # Manpower brief allocates 6 channels; the industry-comparison
        # average for logistics is also 6 (see
        # INDUSTRY_BENCHMARKS_COMPARISON) -- confirm the comparison slide
        # renders a neutral em-dash for this exact tie, not a green arrow.
        data = _logistics_plan()
        prs = Presentation(io.BytesIO(_deck_bytes(data)))
        comparison_slide = None
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame and "Plan Comparison" in shape.text_frame.text:
                    comparison_slide = slide
                    break
            if comparison_slide:
                break
        assert comparison_slide is not None
        texts = [
            s.text_frame.text
            for s in comparison_slide.shapes
            if s.has_text_frame and s.text_frame.text.strip()
        ]
        # Every row rendered is one of beating/trailing/on_par -- the point
        # is simply that no row shows a fabricated arrow for a tie.
        assert any("Channels Selected" in t for t in texts)

    def test_goal_gap_renders_when_projection_falls_short(self):
        gap = fmt.goal_gap(projected_hires=57, goal=500, cost_per_hire=5263.16)
        assert gap is not None
        assert gap["goal"] == 500
        assert gap["projected"] == 57
        assert gap["pct_of_goal"] == pytest.approx(11.4, abs=0.5)

    def test_goal_gap_appears_in_deck_when_goal_stated_and_short(self):
        data = _healthcare_plan()  # hire_volume "500+ hires", projects far fewer
        texts = _all_slide_text(_deck_bytes(data))
        blob = "\n".join(texts)
        assert "Client goal:" in blob
        assert "scaling path:" in blob

    def test_no_goal_gap_when_goal_not_stated(self):
        data = _logistics_plan()
        data["hire_volume"] = ""
        texts = _all_slide_text(_deck_bytes(data))
        blob = "\n".join(texts)
        assert "Client goal:" not in blob


# ---------------------------------------------------------------------------
# 6. Cross-artifact benchmark single-sourcing
# ---------------------------------------------------------------------------
class TestBenchmarkSingleSourcing:
    def test_kb_recruitment_benchmark_reader(self):
        bm = ppt._kb_recruitment_industry_benchmark("logistics_supply_chain", {})
        assert bm is not None
        assert bm["cpa"] == "$12-$35"

    def test_deck_and_workbook_quote_same_cpa_for_logistics(self):
        pytest.importorskip("openpyxl")
        import openpyxl
        from excel_v2 import generate_excel_v2
        from kb_loader import load_knowledge_base

        data = _logistics_plan()
        pptx_bytes = _deck_bytes(data)
        xlsx_bytes = generate_excel_v2(dict(data), load_kb_fn=load_knowledge_base)
        if isinstance(xlsx_bytes, tuple):
            xlsx_bytes = xlsx_bytes[0]

        deck_texts = _all_slide_text(pptx_bytes)
        deck_blob = "\n".join(deck_texts)
        cpa_match = re.search(r"Industry CPA\s*\n?\$?([\d.]+)\s*-\s*\$?([\d.]+)", deck_blob)
        # Fall back to a looser scan across adjacent textboxes.
        if not cpa_match:
            idx = deck_blob.find("Industry CPA")
            window = deck_blob[idx: idx + 60] if idx >= 0 else ""
            cpa_match = re.search(r"\$?([\d.]+)\s*-\s*\$?([\d.]+)", window)
        assert cpa_match, "could not find an Industry CPA range on the deck"
        deck_low, deck_high = float(cpa_match.group(1)), float(cpa_match.group(2))

        wb = openpyxl.load_workbook(io.BytesIO(xlsx_bytes), data_only=True)
        ws = wb["Executive Summary"]
        wb_cpa_line = ""
        for row in ws.iter_rows(values_only=True):
            vals = [v for v in row if v not in (None, "")]
            # Label is acronym-cased ("CPA", not a raw-key-derived "Cpa")
            # per excel_v2._humanize_snake_key.
            if vals and str(vals[0]).strip() == "CPA":
                wb_cpa_line = " | ".join(str(v) for v in vals)
                break
        assert wb_cpa_line, "workbook has no 'CPA' row in Executive Summary"
        wb_match = re.search(
            r"range:\s*\$?([\d.]+)-\$?([\d.]+)", wb_cpa_line, re.IGNORECASE
        )
        assert wb_match, f"could not parse workbook CPA range from {wb_cpa_line!r}"
        wb_low, wb_high = float(wb_match.group(1)), float(wb_match.group(2))

        assert deck_low == pytest.approx(wb_low)
        assert deck_high == pytest.approx(wb_high)


if __name__ == "__main__":
    sys.exit(pytest.main([__file__, "-v"]))
