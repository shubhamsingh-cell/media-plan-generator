"""ppt_generator.BENCHMARKS fallback-dict CPC citation honesty (2026-08-04).

BENCHMARKS is layer 3 (last resort) of _get_benchmarks' 4-layer cascade
(live ad_platform_analysis > KB industry_benchmarks > trend_engine >
BENCHMARKS static dict). Its "cpc" figure in all 9 industries was an
uncited hand-typed literal -- e.g. healthcare_medical and finance_banking
both carried the identical "$0.90 - $3.50" despite having no relationship,
and none of the 9 matched the cited KB figure for that industry.

Fixed by reconciling each "cpc" value to
data/recruitment_benchmarks_deep.json's recruitment_benchmarks
.industry_benchmarks[industry]["cpc"]["range"] field (Appcast/Recruitics/
SHRM-cited, the SAME section _get_benchmarks' Layer 1.5
(_kb_recruitment_industry_benchmark) reads to override this fallback when
available) -- dict structure and the cpa/cph/apply_rate fields (a
different, out-of-scope metric) are unchanged.

Runs under pytest, or standalone:
``python3 tests/test_ppt_benchmarks_fallback_honesty.py``.
"""

from __future__ import annotations

import sys
from pathlib import Path

PROJECT_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(PROJECT_ROOT))

import ppt_generator as ppt  # noqa: E402
from kb_loader import load_knowledge_base  # noqa: E402

# This dict's key -> the KB's real industry_benchmarks key. Two differ
# (tech_engineering/blue_collar_trades) -- a pre-existing, unrelated
# key-mismatch bug in _kb_recruitment_industry_benchmark's runtime lookup,
# not fixed here; the cited figure was still looked up under the KB's real
# key and reconciled into this dict.
_BENCHMARKS_KEY_TO_KB_KEY = {
    "healthcare_medical": "healthcare_medical",
    "tech_engineering": "technology_engineering",
    "retail_consumer": "retail_consumer",
    "finance_banking": "finance_banking",
    "logistics_supply_chain": "logistics_supply_chain",
    "hospitality_travel": "hospitality_travel",
    "blue_collar_trades": "blue_collar_skilled_trades",
    "pharma_biotech": "pharma_biotech",
}

# Retired uncited bands, keyed by the BENCHMARKS industry that carried them.
_RETIRED_UNCITED_CPC = {
    "healthcare_medical": "$0.90 - $3.50",
    "tech_engineering": "$1.20 - $4.50",
    "retail_consumer": "$0.25 - $1.00",
    "finance_banking": "$0.90 - $3.50",
    "logistics_supply_chain": "$0.40 - $1.80",
    "hospitality_travel": "$0.22 - $1.00",
    "blue_collar_trades": "$0.40 - $1.60",
    "pharma_biotech": "$1.50 - $5.00",
}


def _kb_industry_benchmarks() -> dict:
    kb = load_knowledge_base()
    return kb.get("recruitment_benchmarks", {}).get("industry_benchmarks", {})


def test_reconciled_industries_trace_verbatim_to_the_cited_kb_range():
    """Every reconciled industry's cpc figure must equal the KB's own
    range string exactly (not just contain a matching dollar token) --
    this is a straight reconciliation, not a derived/rounded figure."""
    ind_bm = _kb_industry_benchmarks()
    for bm_key, kb_key in _BENCHMARKS_KEY_TO_KB_KEY.items():
        cited_range = ind_bm[kb_key]["cpc"]["range"]
        assert ppt.BENCHMARKS[bm_key]["cpc"] == cited_range, (
            f"BENCHMARKS[{bm_key!r}]['cpc'] = {ppt.BENCHMARKS[bm_key]['cpc']!r} "
            f"does not match cited KB {kb_key}.cpc.range = {cited_range!r}"
        )


def test_retired_uncited_bands_are_gone():
    for bm_key, retired in _RETIRED_UNCITED_CPC.items():
        assert ppt.BENCHMARKS[bm_key]["cpc"] != retired, (
            f"BENCHMARKS[{bm_key!r}]['cpc'] still carries the retired "
            f"uncited band {retired!r}"
        )


def test_general_entry_level_kept_not_invented():
    """general_entry_level has no KB industry_benchmarks counterpart --
    confirm that's still true (sanity on the "no citable figure" premise)
    and that its cpc figure was deliberately left as the pre-existing
    estimate, not replaced with a fabricated "blended" number."""
    ind_bm = _kb_industry_benchmarks()
    assert "general_entry_level" not in ind_bm
    assert ppt.BENCHMARKS["general_entry_level"]["cpc"] == "$0.35 - $1.30"


def test_cpa_cph_apply_rate_untouched():
    """Only cpc changed -- cpa/cph/apply_rate are a different metric and
    were out of scope (not flagged CPC_UNCITED_SURFACE)."""
    expected_non_cpc = {
        "healthcare_medical": {"cpa": "$35 - $85", "cph": "$9K - $12K", "apply_rate": "3.2% - 4.5%"},
        "tech_engineering": {"cpa": "$25 - $75", "cph": "$6K - $22K", "apply_rate": "6.41%"},
        "retail_consumer": {"cpa": "$8 - $21", "cph": "$2.7K - $4K", "apply_rate": "4.5% - 5.8%"},
        "general_entry_level": {"cpa": "$10 - $25", "cph": "$2K - $4.7K", "apply_rate": "5.5% - 6.1%"},
        "finance_banking": {"cpa": "$21 - $65", "cph": "$5K - $12K", "apply_rate": "5.0% - 6.0%"},
        "logistics_supply_chain": {"cpa": "$15 - $52", "cph": "$4.5K - $8K", "apply_rate": "4.0% - 5.2%"},
        "hospitality_travel": {"cpa": "$8 - $25", "cph": "$2.5K - $4K", "apply_rate": "4.0% - 5.0%"},
        "blue_collar_trades": {"cpa": "$12 - $35", "cph": "$3.5K - $5.6K", "apply_rate": "4.0% - 5.5%"},
        "pharma_biotech": {"cpa": "$40 - $110", "cph": "$8K - $18K", "apply_rate": "3.8% - 5.2%"},
    }
    for industry, fields in expected_non_cpc.items():
        for field, value in fields.items():
            assert ppt.BENCHMARKS[industry][field] == value, (
                f"BENCHMARKS[{industry!r}][{field!r}] changed -- out of scope for this fix"
            )


def test_fallback_dict_is_irrelevant_when_kb_layer_1_5_is_active(monkeypatch):
    """Fallback-only proof: for an industry an upper cascade layer DOES cover,
    _get_benchmarks' result must be identical regardless of what
    BENCHMARKS[industry]['cpc'] holds. Proven by swapping in the OLD retired
    uncited value at runtime and showing the final cpc doesn't move.

    6 of the 8 reconciled industries are covered by Layer 1.5
    (_kb_recruitment_industry_benchmark), which unconditionally overrides the
    fallback and sets confidence="market_intelligence_kb" -- asserted
    strictly for those. The other 2 (tech_engineering, blue_collar_trades)
    hit the pre-existing, out-of-scope key-mismatch bug documented in the
    module comment above (this dict's key doesn't match the KB's key), so
    Layer 1.5 silently no-ops for them at runtime and Layer 2 (trend_engine)
    is the layer actually shadowing BENCHMARKS instead -- confidence for
    those two is asserted to be "trend_engine", not "market_intelligence_kb",
    and NOT "curated" (which would mean nothing shadowed BENCHMARKS at all).
    """
    kb = load_knowledge_base()
    data = {"_knowledge_base": kb}

    for bm_key, retired_old_value in _RETIRED_UNCITED_CPC.items():
        kb_key_matches = _BENCHMARKS_KEY_TO_KB_KEY[bm_key] == bm_key

        result_with_new_fallback = ppt._get_benchmarks(bm_key, data)

        monkeypatch.setitem(
            ppt.BENCHMARKS, bm_key, {**ppt.BENCHMARKS[bm_key], "cpc": retired_old_value}
        )
        result_with_old_fallback = ppt._get_benchmarks(bm_key, data)
        monkeypatch.undo()  # restore BENCHMARKS[bm_key] before next iteration

        assert result_with_new_fallback["cpc"] == result_with_old_fallback["cpc"], (
            f"{bm_key}: _get_benchmarks' final cpc changed depending on the "
            "layer-3 fallback dict's value even though an upper cascade "
            "layer covers this industry -- the fallback should be fully "
            "shadowed"
        )
        if kb_key_matches:
            assert result_with_new_fallback["confidence"] == "market_intelligence_kb", (
                f"{bm_key}: expected Layer 1.5 (KB) to shadow the fallback"
            )
        else:
            # Known key-mismatch bug (see module comment): Layer 1.5 can't
            # find this industry under its own key, so it no-ops. Layer 2
            # (trend_engine) is the one actually doing the shadowing here --
            # confirm that's really what happened, not a silent fall-through
            # to the static BENCHMARKS dict this test is trying to prove is
            # irrelevant.
            assert result_with_new_fallback["confidence"] == "trend_engine", (
                f"{bm_key}: expected the known KB key-mismatch bug to route "
                "this industry to Layer 2 (trend_engine) instead of Layer "
                f"1.5 -- got confidence={result_with_new_fallback['confidence']!r}"
            )


def test_rendered_channel_strategy_slide_unaffected_when_kb_layer_active():
    """End-to-end byte-level proof: the actual deck slide text for an
    industry the KB covers is identical whether BENCHMARKS carries the new
    cited value or the old retired one -- confirming this fix only changes
    output for plans that fall all the way through to the static fallback
    (no KB, no trend_engine, no live ad_platform_analysis)."""
    from pptx import Presentation

    import budget_engine

    kb = load_knowledge_base()
    industry = "healthcare_medical"
    data = {
        "client_name": "Acme",
        "industry": industry,
        "locations": ["United States"],
        "roles": ["Registered Nurse"],
        "_knowledge_base": kb,
        "_budget_allocation": budget_engine.calculate_budget_allocation(
            total_budget=100_000,
            roles=[{"title": "Registered Nurse", "count": 20, "tier": "mid"}],
            locations=[{"city": "Metro", "state": "", "country": "United States"}],
            industry=industry,
            channel_percentages={"Indeed": 50, "LinkedIn": 50},
            collar_type="white",
            campaign_start_month=9,
        ),
    }

    def _render_texts():
        ppt._set_active_currency(data)
        prs = Presentation()
        prs.slide_width = ppt.SLIDE_WIDTH
        prs.slide_height = ppt.SLIDE_HEIGHT
        ppt._build_slide_channel_strategy(prs, data)
        texts = [
            shape.text_frame.text
            for shape in prs.slides[0].shapes
            if shape.has_text_frame
        ]
        ppt._set_active_currency({})
        return texts

    texts_with_new_fallback = _render_texts()

    old_cpc = ppt.BENCHMARKS[industry]["cpc"]
    ppt.BENCHMARKS[industry]["cpc"] = "$0.90 - $3.50"  # the retired uncited band
    try:
        texts_with_old_fallback = _render_texts()
    finally:
        ppt.BENCHMARKS[industry]["cpc"] = old_cpc

    assert texts_with_new_fallback == texts_with_old_fallback, (
        "Channel Strategy slide text differs depending on the layer-3 "
        "fallback dict even though this plan's industry is KB-covered -- "
        "the fallback dict must be fully shadowed by Layer 1.5"
    )


if __name__ == "__main__":
    import pytest

    raise SystemExit(pytest.main([__file__, "-v"]))
