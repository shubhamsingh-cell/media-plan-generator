"""skill_target.py follow-up to 8fc9a4f: `X or [][:N]` operator-precedence
no-ops, verified by targeted repro before fixing.

Python precedence: slicing binds tighter than `or`, so
``analysis.get("occupations") or [][:20]`` slices the empty-list LITERAL --
the intended cap is silently a no-op and the real list flows through
unbounded. Ten sites in this module carried the pattern:

  1.  skill_target.py:1730  recommend_channels -- COLLAR_STRATEGY[collar]
      "preferred_platforms" or [][:3] (top-3 platform boost eligibility)
  2.  skill_target.py:2130  generate_skill_excel -- occupations or [][:20]
      (Matched Occupations sheet rows)
  3.  skill_target.py:2228  generate_skill_excel -- hotspots or [][:15]
      (Hotspots sheet rows)
  4.  skill_target.py:2251  generate_skill_excel -- adjacent_skills or [][:20]
      (Adjacent Skills sheet rows)
  5.  skill_target.py:2387  generate_skill_ppt -- occupations or [][:8]
      (Slide 2 table rows)
  6.  skill_target.py:2451  generate_skill_ppt -- channels or [][:8]
      (Slide 4 table rows)
  7.  skill_target.py:2491  generate_skill_ppt -- hotspots or [][:5]
      (Slide 5 table rows)
  8.  skill_target.py:2518  generate_skill_ppt -- top_employers or [][:4]
      (Slide 5 hotspots table, Top Employers cell)
  9.  skill_target.py:2536  generate_skill_ppt -- adjacent_skills or [][:8]
      (Slide 6 table rows)
  10. skill_target.py:2548  generate_skill_ppt -- suggested_titles or [][:10]
      (Slide 6 job-titles textbox)

Each site fans out into a rendered artifact (an .xlsx sheet, a .pptx table
or textbox, or a channel-weight boost) -- restoring the cap SHRINKS that
output, which is the correct behaviour: the PPT tables already clamp their
pixel HEIGHT via `min((len(x) + 1) * 0.4, cap)` while the row COUNT was
unbounded, so an uncapped list previously squeezed arbitrarily many rows
into a fixed-height table.

A source-level guard test keeps the pattern from reappearing in this file.

Runs under pytest, or standalone:
``python3 tests/test_truncation_precedence_skill_target.py``.
"""

from __future__ import annotations

import io
import re
import sys
from pathlib import Path

PROJECT_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(PROJECT_ROOT))

import openpyxl  # noqa: E402
from pptx import Presentation  # noqa: E402

import skill_target  # noqa: E402


# ===========================================================================
# Guard: the precedence pattern must not reappear in this file
# ===========================================================================
def test_slice_precedence_pattern_absent_from_skill_target():
    """`or ""[...]` / `or ''[...]` / `or [][...]` always slice the empty
    LITERAL -- any occurrence is a bug by construction."""
    pattern = re.compile(r"""or\s+(""|''|\[\])\s*\[""")
    src = (PROJECT_ROOT / "skill_target.py").read_text(encoding="utf-8")
    hits = [
        f"skill_target.py:{i}"
        for i, line in enumerate(src.splitlines(), 1)
        if pattern.search(line)
    ]
    assert not hits, f"slice-binds-to-literal precedence bug at: {hits}"


# ---------------------------------------------------------------------------
# Shared stub-data builders
# ---------------------------------------------------------------------------
def _make_occ(i: int) -> dict:
    return {
        "soc": f"SOC{i}",
        "title": f"Occupation {i}",
        "zone": "3",
        "match_count": 1,
        "matched_skills": ["python"],
    }


def _make_hotspot(i: int, employers=None) -> dict:
    return {
        "metro": f"Metro {i}",
        "avg_concentration": 0.5,
        "skill_coverage": 0.5,
        "top_employers": (
            employers if employers is not None else [f"Employer{j}" for j in range(8)]
        ),
    }


def _make_adj(i: int) -> dict:
    return {
        "skill": f"AdjSkill{i}",
        "max_relevance": 0.5,
        "scarcity": 0.5,
        "connected_to": ["python"],
    }


def _base_analysis(**overrides) -> dict:
    data = {
        "input": {"skills": ["python"], "industry": "technology", "location": "Dallas, TX"},
        "summary": "",
        "occupations": [],
        "demand_trends": {"summary": {}, "skills": {}},
        "channels": [],
        "salary_benchmarks": {"occupations": {}},
        "hotspots": [],
        "adjacent_skills": [],
        "suggested_titles": [],
    }
    data.update(overrides)
    return data


def _sheet_data_row_count(ws, min_row: int) -> int:
    """Count rows with a non-empty value in column B (the first data col)."""
    return sum(1 for r in ws.iter_rows(min_row=min_row, min_col=2, max_col=2) if r[0].value)


def _tables_on_slide(slide):
    return [sh.table for sh in slide.shapes if sh.has_table]


# ===========================================================================
# Sites 2, 3, 4: generate_skill_excel (Matched Occupations / Hotspots /
# Adjacent Skills sheet row counts)
# ===========================================================================
def test_excel_occupations_capped_at_twenty():
    analysis = _base_analysis(occupations=[_make_occ(i) for i in range(25)])
    wb = openpyxl.load_workbook(io.BytesIO(skill_target.generate_skill_excel(analysis)))
    assert _sheet_data_row_count(wb["Matched Occupations"], min_row=6) == 20


def test_excel_occupations_short_input_unchanged():
    analysis = _base_analysis(occupations=[_make_occ(i) for i in range(5)])
    wb = openpyxl.load_workbook(io.BytesIO(skill_target.generate_skill_excel(analysis)))
    assert _sheet_data_row_count(wb["Matched Occupations"], min_row=6) == 5


def test_excel_hotspots_capped_at_fifteen():
    analysis = _base_analysis(hotspots=[_make_hotspot(i) for i in range(25)])
    wb = openpyxl.load_workbook(io.BytesIO(skill_target.generate_skill_excel(analysis)))
    assert _sheet_data_row_count(wb["Hotspots"], min_row=4) == 15


def test_excel_hotspots_short_input_unchanged():
    analysis = _base_analysis(hotspots=[_make_hotspot(i) for i in range(3)])
    wb = openpyxl.load_workbook(io.BytesIO(skill_target.generate_skill_excel(analysis)))
    assert _sheet_data_row_count(wb["Hotspots"], min_row=4) == 3


def test_excel_adjacent_skills_capped_at_twenty():
    analysis = _base_analysis(adjacent_skills=[_make_adj(i) for i in range(25)])
    wb = openpyxl.load_workbook(io.BytesIO(skill_target.generate_skill_excel(analysis)))
    assert _sheet_data_row_count(wb["Adjacent Skills"], min_row=4) == 20


def test_excel_adjacent_skills_short_input_unchanged():
    analysis = _base_analysis(adjacent_skills=[_make_adj(i) for i in range(4)])
    wb = openpyxl.load_workbook(io.BytesIO(skill_target.generate_skill_excel(analysis)))
    assert _sheet_data_row_count(wb["Adjacent Skills"], min_row=4) == 4


def test_excel_none_lists_render_empty_without_crash():
    """`None` (not just a missing key or `[]`) must flow through the `or []`
    guard cleanly -- the workbook must still build with zero data rows."""
    analysis = _base_analysis(occupations=None, hotspots=None, adjacent_skills=None)
    xlsx_bytes = skill_target.generate_skill_excel(analysis)
    assert xlsx_bytes, "generator must not silently fail (returns b'' on exception)"
    wb = openpyxl.load_workbook(io.BytesIO(xlsx_bytes))
    assert _sheet_data_row_count(wb["Matched Occupations"], min_row=6) == 0
    assert _sheet_data_row_count(wb["Hotspots"], min_row=4) == 0
    assert _sheet_data_row_count(wb["Adjacent Skills"], min_row=4) == 0


# ===========================================================================
# Sites 5, 6, 9, 10: generate_skill_ppt table / textbox item counts
# ===========================================================================
def test_ppt_occupations_capped_at_eight():
    analysis = _base_analysis(occupations=[_make_occ(i) for i in range(25)])
    prs = Presentation(io.BytesIO(skill_target.generate_skill_ppt(analysis)))
    table = _tables_on_slide(prs.slides[1])[0]
    assert len(table.rows) == 9  # 8 data rows + 1 header


def test_ppt_occupations_short_input_unchanged():
    analysis = _base_analysis(occupations=[_make_occ(i) for i in range(3)])
    prs = Presentation(io.BytesIO(skill_target.generate_skill_ppt(analysis)))
    table = _tables_on_slide(prs.slides[1])[0]
    assert len(table.rows) == 4


def test_ppt_channels_capped_at_eight():
    analysis = _base_analysis(
        channels=[{"name": f"Ch{i}", "weight": 0.1, "reasons": ["r"]} for i in range(25)]
    )
    prs = Presentation(io.BytesIO(skill_target.generate_skill_ppt(analysis)))
    table = _tables_on_slide(prs.slides[3])[0]
    assert len(table.rows) == 9


def test_ppt_channels_short_input_unchanged():
    analysis = _base_analysis(
        channels=[{"name": f"Ch{i}", "weight": 0.1, "reasons": ["r"]} for i in range(2)]
    )
    prs = Presentation(io.BytesIO(skill_target.generate_skill_ppt(analysis)))
    table = _tables_on_slide(prs.slides[3])[0]
    assert len(table.rows) == 3


def test_ppt_adjacent_skills_capped_at_eight():
    analysis = _base_analysis(adjacent_skills=[_make_adj(i) for i in range(25)])
    prs = Presentation(io.BytesIO(skill_target.generate_skill_ppt(analysis)))
    table = _tables_on_slide(prs.slides[5])[0]
    assert len(table.rows) == 9


def test_ppt_adjacent_skills_short_input_unchanged():
    analysis = _base_analysis(adjacent_skills=[_make_adj(i) for i in range(3)])
    prs = Presentation(io.BytesIO(skill_target.generate_skill_ppt(analysis)))
    table = _tables_on_slide(prs.slides[5])[0]
    assert len(table.rows) == 4


def _titles_textbox_text(slide):
    return next(
        sh.text_frame.text
        for sh in slide.shapes
        if sh.has_text_frame and "Title 0" in sh.text_frame.text
    )


def test_ppt_suggested_titles_capped_at_ten():
    analysis = _base_analysis(suggested_titles=[f"Title {i}" for i in range(25)])
    prs = Presentation(io.BytesIO(skill_target.generate_skill_ppt(analysis)))
    text = _titles_textbox_text(prs.slides[5])
    assert text.count("Title ") == 10
    assert "Title 10" not in text  # the 11th item (index 10) must be dropped


def test_ppt_suggested_titles_short_input_unchanged():
    analysis = _base_analysis(suggested_titles=[f"Title {i}" for i in range(3)])
    prs = Presentation(io.BytesIO(skill_target.generate_skill_ppt(analysis)))
    text = _titles_textbox_text(prs.slides[5])
    assert text.count("Title ") == 3


# ===========================================================================
# Sites 7, 8: generate_skill_ppt Hotspots table (row cap 5) + Top Employers
# cell (item cap 4) -- share a table, tested together.
# ===========================================================================
def test_ppt_hotspots_capped_at_five():
    analysis = _base_analysis(hotspots=[_make_hotspot(i) for i in range(25)])
    prs = Presentation(io.BytesIO(skill_target.generate_skill_ppt(analysis)))
    hs_table = _tables_on_slide(prs.slides[4])[-1]
    assert len(hs_table.rows) == 6  # 5 data rows + 1 header


def test_ppt_hotspots_short_input_unchanged():
    analysis = _base_analysis(hotspots=[_make_hotspot(i) for i in range(2)])
    prs = Presentation(io.BytesIO(skill_target.generate_skill_ppt(analysis)))
    hs_table = _tables_on_slide(prs.slides[4])[-1]
    assert len(hs_table.rows) == 3


def test_ppt_top_employers_capped_at_four():
    analysis = _base_analysis(
        hotspots=[_make_hotspot(0, employers=[f"Employer{j}" for j in range(8)])]
    )
    prs = Presentation(io.BytesIO(skill_target.generate_skill_ppt(analysis)))
    hs_table = _tables_on_slide(prs.slides[4])[-1]
    cell_text = hs_table.cell(1, 3).text
    assert cell_text.count(",") + 1 == 4
    assert "Employer4" not in cell_text


def test_ppt_top_employers_short_input_unchanged():
    analysis = _base_analysis(hotspots=[_make_hotspot(0, employers=["A", "B"])])
    prs = Presentation(io.BytesIO(skill_target.generate_skill_ppt(analysis)))
    hs_table = _tables_on_slide(prs.slides[4])[-1]
    assert hs_table.cell(1, 3).text == "A, B"


def test_ppt_none_lists_render_without_crash():
    """`None` for every capped list must not raise -- each section is
    conditionally rendered only `if <list>:`, so a `None` -> `[]` result
    must simply skip that table/textbox."""
    analysis = _base_analysis(
        occupations=None,
        channels=None,
        hotspots=None,
        adjacent_skills=None,
        suggested_titles=None,
    )
    ppt_bytes = skill_target.generate_skill_ppt(analysis)
    assert ppt_bytes, "generator must not silently fail (returns b'' on exception)"
    prs = Presentation(io.BytesIO(ppt_bytes))
    assert len(prs.slides) == 6
    assert _tables_on_slide(prs.slides[1]) == []
    assert _tables_on_slide(prs.slides[3]) == []
    assert _tables_on_slide(prs.slides[4]) == []
    assert _tables_on_slide(prs.slides[5]) == []


# ===========================================================================
# Site 1: recommend_channels -- COLLAR_STRATEGY preferred_platforms cap (3)
# ===========================================================================
def _patch_collar_env(preferred_platforms, channel_names):
    """Swap in a stub COLLAR_STRATEGY + _SKILL_CATEGORY_CHANNELS so the
    preferred-platform boost is observable against known, equal-weight
    channels. Returns the (weights dict) after calling recommend_channels
    with an unmapped skill (forces the 'business' category fallback)."""
    orig_collar_strategy = skill_target.COLLAR_STRATEGY
    orig_has_collar = skill_target._HAS_COLLAR
    orig_channels = skill_target._SKILL_CATEGORY_CHANNELS
    try:
        skill_target.COLLAR_STRATEGY = {
            "test_collar": {"preferred_platforms": preferred_platforms}
        }
        skill_target._HAS_COLLAR = True
        skill_target._SKILL_CATEGORY_CHANNELS = {
            "business": {
                "channels": [
                    {"name": f"{n} Board", "weight": 0.1, "reason": "r"}
                    for n in channel_names
                ]
            }
        }
        result = skill_target.recommend_channels(
            {"skills": ["totally_unmapped_skill_xyz"], "collar_type": "test_collar"}
        )
        return {c["name"]: c["weight"] for c in result}
    finally:
        skill_target.COLLAR_STRATEGY = orig_collar_strategy
        skill_target._HAS_COLLAR = orig_has_collar
        skill_target._SKILL_CATEGORY_CHANNELS = orig_channels


def test_recommend_channels_preferred_platforms_capped_at_three():
    """6 preferred_platforms, 6 matching equal-weight channels: only the
    first 3 should receive the 1.15x boost. If the cap is a no-op, all 6
    channels come out with an identical (uniformly boosted) weight."""
    names = ["Alpha", "Beta", "Gamma", "Delta", "Echo", "Foxtrot"]
    weights = _patch_collar_env(preferred_platforms=names, channel_names=names)
    boosted = {weights["Alpha Board"], weights["Beta Board"], weights["Gamma Board"]}
    unboosted = {weights["Delta Board"], weights["Echo Board"], weights["Foxtrot Board"]}
    assert len(boosted) == 1, "the first 3 preferred platforms must share one weight"
    assert len(unboosted) == 1, "the last 3 (uncapped-in) platforms must share one weight"
    assert next(iter(boosted)) > next(iter(unboosted)), (
        "only the first 3 preferred_platforms should get the 1.15x boost -- "
        "if all 6 are boosted equally, the [:3] cap is a no-op"
    )


def test_recommend_channels_short_platform_list_unchanged():
    """<= 3 preferred_platforms: the cap changes nothing, both must still
    be boosted relative to an unmatched channel."""
    weights = _patch_collar_env(
        preferred_platforms=["Alpha", "Beta"],
        channel_names=["Alpha", "Beta", "Gamma"],
    )
    assert weights["Alpha Board"] == weights["Beta Board"]
    assert weights["Alpha Board"] > weights["Gamma Board"]


def test_recommend_channels_missing_preferred_platforms_no_crash():
    """A collar entry with no `preferred_platforms` key at all (None via
    `.get`) must not raise -- the boost loop simply doesn't run."""
    orig_collar_strategy = skill_target.COLLAR_STRATEGY
    orig_has_collar = skill_target._HAS_COLLAR
    try:
        skill_target.COLLAR_STRATEGY = {"test_collar": {}}
        skill_target._HAS_COLLAR = True
        result = skill_target.recommend_channels(
            {"skills": ["python"], "collar_type": "test_collar"}
        )
        assert isinstance(result, list) and result
    finally:
        skill_target.COLLAR_STRATEGY = orig_collar_strategy
        skill_target._HAS_COLLAR = orig_has_collar


if __name__ == "__main__":
    _failures = 0
    for _name, _fn in sorted(globals().items()):
        if _name.startswith("test_") and callable(_fn):
            try:
                _fn()
                print(f"PASS {_name}")
            except AssertionError as exc:
                _failures += 1
                print(f"FAIL {_name}: {exc}")
            except Exception as exc:  # noqa: BLE001
                _failures += 1
                print(f"ERROR {_name}: {exc}")
    if _failures:
        sys.exit(1)
