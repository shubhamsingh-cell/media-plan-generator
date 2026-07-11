"""Tests for Agent W4B's ppt_generator.py / display_format.py bundle-quality
fixes (rescore3).

Covers:
    1. strategy:manpower#3: the comparison slide's "vs. Client Goal" row is
       rendered as its own distinctly-labeled band, never inside the
       "Industry Average" column (a category error the old layout made).
    2. strategy:manpower#2: two unsourced deck statistics -- "$52+ CPA" and
       the fabricated "10,238+ publishers" default -- are gone; each is
       either sourced from the KB or replaced with a non-numeric claim.
    3. copy:both#2: Implementation Timeline phase action items interpolate
       this plan's own facts (client name, top channels) instead of being
       byte-identical across industries/durations.
    4. copy:both#5-family: acronym-preserving title case ("CDL Drivers",
       not "Cdl Drivers") in display_format.smart_title, used at the deck
       call site that produced the bug.

Runs under pytest, or standalone: ``python3 tests/test_ppt_w4b_bundle_quality.py``.
"""

from __future__ import annotations

import os
import sys

import pytest

PROJECT_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, PROJECT_ROOT)

import ppt_generator as ppt  # noqa: E402
import display_format as fmt  # noqa: E402
from pptx import Presentation  # noqa: E402


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def _new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = ppt.SLIDE_WIDTH
    prs.slide_height = ppt.SLIDE_HEIGHT
    return prs


def _all_slide_text(prs: Presentation) -> list[str]:
    out: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                out.append(shape.text_frame.text)
    return out


def _base_plan(**overrides) -> dict:
    data = {
        "client_name": "Manpower - Amerigas",
        "industry": "logistics_supply_chain",
        "industry_label": "Logistics & Supply Chain",
        "budget": "$150,000",
        "campaign_duration": "6 months",
        "campaign_weeks": 24,
        "locations": ["Massachusetts", "Maine", "New Hampshire"],
        "roles": ["CDL A Driver"],
        "target_roles": [{"title": "CDL A Driver", "count": 100, "tier": "Hourly"}],
        "work_environment": "on_site",
        "hire_volume": "100 hires",
        "_budget_allocation": {
            "channel_allocations": {
                "programmatic_dsp": {"dollar_amount": 39_535.5, "percentage": 26.4},
                "regional_boards": {"dollar_amount": 36_175.5, "percentage": 24.1},
                "global_boards": {"dollar_amount": 33_655.5, "percentage": 22.4},
                "niche_boards": {"dollar_amount": 29_221.5, "percentage": 19.5},
                "social_media": {"dollar_amount": 3_912.0, "percentage": 2.6},
                "employer_branding": {"dollar_amount": 7_500.0, "percentage": 5.0},
            },
            "metadata": {"total_budget": 150_000},
            "total_projected": {"hires": 48, "applications": 18_950},
        },
    }
    data.update(overrides)
    return data


# ---------------------------------------------------------------------------
# 1. strategy:manpower#3 -- client goal is never shown as "industry average"
# ---------------------------------------------------------------------------
class TestClientGoalNeverInBenchmarkColumn:
    def test_goal_row_rendered_as_own_labeled_band(self):
        data = _base_plan(hire_volume="100 hires")  # 48 projected < 100 goal
        prs = _new_prs()
        ppt._build_slide_comparison_timeline(prs, data)
        blob = "\n".join(_all_slide_text(prs))
        assert "CLIENT GOAL" in blob
        assert "100 hires target" in blob or "100" in blob

    def test_goal_row_never_labeled_industry_average(self):
        # The old bug rendered "100 hires goal" as the value in the column
        # headed "... Average" -- i.e. a client target dressed up as a
        # benchmark. The goal callout text itself must say it is NOT a
        # benchmark.
        data = _base_plan(hire_volume="100 hires")
        prs = _new_prs()
        ppt._build_slide_comparison_timeline(prs, data)
        blob = "\n".join(_all_slide_text(prs))
        assert "not an industry benchmark" in blob.lower()

    def test_no_goal_gap_means_no_goal_band(self):
        # Projected hires (48) already meets/exceeds a modest goal -- no
        # gap, so no goal band should render at all.
        data = _base_plan(hire_volume="10 hires")
        prs = _new_prs()
        ppt._build_slide_comparison_timeline(prs, data)
        blob = "\n".join(_all_slide_text(prs))
        assert "CLIENT GOAL" not in blob


# ---------------------------------------------------------------------------
# 2. strategy:manpower#2 -- no unsourced/fabricated numeric claims
# ---------------------------------------------------------------------------
class TestNoUnsourcedDeckStatistics:
    def test_complications_no_longer_claim_52_plus_cpa(self):
        complications = ppt._get_complications("logistics_supply_chain")
        blob = "\n".join(complications)
        assert "$52+" not in blob
        # The corrected, KB-sourced CDL benchmark should appear instead.
        assert "$25-$50" in blob

    def test_publisher_count_falls_back_to_non_numeric_claim_without_kb(self):
        # No "_joveo_publishers" and no "_knowledge_base" attached -- the
        # real KB file may still resolve via kb_loader, so this asserts the
        # helper never returns the old hardcoded-wrong 10238 default.
        total = ppt._joveo_total_active_publishers({})
        assert total != 10238

    def test_publisher_count_uses_real_pipeline_value_when_present(self):
        data = {"_joveo_publishers": {"total_active_publishers": 1238}}
        assert ppt._joveo_total_active_publishers(data) == 1238

    def test_cpc_trend_not_leaked_when_kb_lacks_trend_yoy(self):
        # logistics_supply_chain's KB cpc node has no "trend_yoy" field --
        # any cpc_trend must not silently carry trend_engine's independent
        # (unsourced-in-workbook) estimate.
        kb = {
            "recruitment_benchmarks": {
                "industry_benchmarks": {
                    "logistics_supply_chain": {
                        "cpa": {"range": "$12-$35", "trend_yoy": "+8-12%"},
                        "cpc": {"range": "$0.30-$1.50"},  # no trend_yoy key
                    }
                }
            }
        }
        data = {"_knowledge_base": kb, "industry": "logistics_supply_chain"}
        bm = ppt._get_benchmarks("logistics_supply_chain", data)
        assert not bm.get("cpc_trend")


# ---------------------------------------------------------------------------
# 3. copy:both#2 -- Implementation Timeline bullets are plan-specific
# ---------------------------------------------------------------------------
class TestTimelineBulletsInterpolated:
    def test_kickoff_bullet_names_the_client(self):
        data = _base_plan(client_name="Manpower - Amerigas")
        prs = _new_prs()
        ppt._build_slide_comparison_timeline(prs, data)
        blob = "\n".join(_all_slide_text(prs))
        assert "Manpower - Amerigas" in blob

    def test_scale_bullet_names_top_channels(self):
        data = _base_plan()
        prs = _new_prs()
        ppt._build_slide_comparison_timeline(prs, data)
        blob = "\n".join(_all_slide_text(prs))
        # Top channel by dollar amount is Programmatic DSP.
        assert "Scale top performers: Programmatic" in blob

    def test_timeline_bullets_differ_across_clients(self):
        manpower = _base_plan(
            client_name="Manpower - Amerigas",
            industry="logistics_supply_chain",
            campaign_weeks=24,
        )
        atria = _base_plan(
            client_name="Atria Senior Living",
            industry="healthcare_medical",
            campaign_weeks=78,
            budget="$300,000",
            _budget_allocation={
                "channel_allocations": {
                    "programmatic_dsp": {"dollar_amount": 84_336.0, "percentage": 28.1},
                    "global_boards": {"dollar_amount": 79_400.0, "percentage": 26.5},
                },
                "metadata": {"total_budget": 300_000},
                "total_projected": {"hires": 57, "applications": 17_780},
            },
        )
        prs1, prs2 = _new_prs(), _new_prs()
        ppt._build_slide_comparison_timeline(prs1, manpower)
        ppt._build_slide_comparison_timeline(prs2, atria)
        blob1 = "\n".join(_all_slide_text(prs1))
        blob2 = "\n".join(_all_slide_text(prs2))
        assert "Manpower - Amerigas" in blob1
        assert "Atria Senior Living" in blob2
        assert "Manpower - Amerigas" not in blob2
        assert "Atria Senior Living" not in blob1


# ---------------------------------------------------------------------------
# 4. copy:both#5-family -- acronym-preserving title case
# ---------------------------------------------------------------------------
class TestSmartTitleAcronyms:
    def test_cdl_drivers_not_clobbered(self):
        assert fmt.smart_title("cdl_drivers") == "CDL Drivers"

    def test_plain_words_still_title_cased(self):
        assert fmt.smart_title("warehouse_hourly") == "Warehouse Hourly"

    def test_mixed_acronym_and_plain_words(self):
        assert fmt.smart_title("rn_and_lpn_roles") == "RN And LPN Roles"

    def test_handles_space_separated_input_too(self):
        assert fmt.smart_title("hvac technician") == "HVAC Technician"

    def test_empty_and_none_safe(self):
        assert fmt.smart_title("") == ""
        assert fmt.smart_title(None) == ""


if __name__ == "__main__":
    sys.exit(pytest.main([__file__, "-v"]))
