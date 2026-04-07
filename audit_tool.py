"""
audit_tool.py -- Recruitment Advertising Media Plan Audit Tool

Parses a client's existing media plan (Excel/CSV), audits every line item
against benchmark data from trend_engine, identifies overspend, flags
missing channels via collar_intelligence + budget_engine, and generates
a savings/improvement report with branded Excel & PPT deliverables.

This is NOT a performance tracker (see performance_tracker.py).
That tool evaluates ACTUAL campaign results. This tool evaluates a
PLANNED media plan before it runs -- catching overspend and coverage
gaps before a single dollar is spent.

Thread-safe, never crashes (all errors return structured error dicts).

Design tokens:
    Excel: Sapphire Blue palette (Navy #0F172A, Sapphire #2563EB, Light #DBEAFE)
    PPT:   Brand (Port Gore #202058, Blue Violet #5A54BD, Downy #6BB3CD)
"""

from __future__ import annotations

import io
import logging
import os
import re
import datetime
from typing import Any, Dict, List, Optional

logger = logging.getLogger(__name__)

# ── Optional imports (lazy, with try/except like performance_tracker.py) ───
try:
    import trend_engine as _trend_engine

    _HAS_TREND_ENGINE = True
except ImportError:
    _trend_engine = None
    _HAS_TREND_ENGINE = False

try:
    import budget_engine as _budget_engine

    _HAS_BUDGET_ENGINE = True
except ImportError:
    _budget_engine = None
    _HAS_BUDGET_ENGINE = False

try:
    import collar_intelligence as _collar_intel

    _HAS_COLLAR_INTEL = True
except ImportError:
    _collar_intel = None
    _HAS_COLLAR_INTEL = False

try:
    from shared_utils import INDUSTRY_LABEL_MAP
except ImportError:
    INDUSTRY_LABEL_MAP = {}

    def parse_budget(v, *, default=100_000.0):
        try:
            return float(v)
        except Exception:
            return default


try:
    from benchmark_registry import get_benchmark_value

    _HAS_BENCHMARK_REGISTRY = True
except ImportError:
    _HAS_BENCHMARK_REGISTRY = False


# ═══════════════════════════════════════════════════════════════════════════════
# COLUMN MATCHING -- flexible, case-insensitive, partial match
# ═══════════════════════════════════════════════════════════════════════════════

_COLUMN_PATTERNS: Dict[str, List[str]] = {
    "channel": [
        "channel",
        "platform",
        "source",
        "medium",
        "publisher",
        "vendor",
        "ad platform",
        "network",
        "job board",
        "site",
        "media",
    ],
    "budget": [
        "budget",
        "spend",
        "cost",
        "investment",
        "total spend",
        "total cost",
        "ad spend",
        "media spend",
        "allocation",
        "planned spend",
        "monthly budget",
        "annual budget",
        "total budget",
    ],
    "cpc": [
        "cpc",
        "cost per click",
        "cost/click",
        "avg cpc",
        "average cpc",
        "planned cpc",
        "target cpc",
        "est cpc",
        "estimated cpc",
    ],
    "cpa": [
        "cpa",
        "cost per application",
        "cost per apply",
        "cost/apply",
        "cost per conversion",
        "cost per lead",
        "cost/application",
        "cost per acquisition",
        "planned cpa",
        "target cpa",
        "est cpa",
    ],
    "target_roles": [
        "role",
        "roles",
        "job title",
        "position",
        "target role",
        "job type",
        "job category",
        "occupation",
    ],
    "target_locations": [
        "location",
        "locations",
        "city",
        "region",
        "geography",
        "geo",
        "target location",
        "market",
        "area",
        "state",
        "country",
    ],
    "duration": [
        "duration",
        "period",
        "timeline",
        "start",
        "end",
        "dates",
        "flight",
        "campaign length",
        "weeks",
        "months",
    ],
    "impressions": [
        "impression",
        "impr",
        "views",
        "estimated impressions",
        "reach",
        "est impressions",
        "projected impressions",
    ],
    "clicks": [
        "click",
        "total click",
        "estimated clicks",
        "projected clicks",
    ],
}


def _match_column(header: str, patterns: List[str]) -> bool:
    """Check if a column header matches any of the given patterns (case-insensitive, partial)."""
    h = header.lower().strip()
    for pat in patterns:
        if pat in h or h in pat:
            return True
    return False


def _map_columns(headers: List[str]) -> Dict[str, Optional[int]]:
    """Map canonical field names to column indices from the header row."""
    mapping: Dict[str, Optional[int]] = {k: None for k in _COLUMN_PATTERNS}
    for idx, header in enumerate(headers):
        if not header or not header.strip():
            continue
        for field, patterns in _COLUMN_PATTERNS.items():
            if mapping[field] is None and _match_column(header, patterns):
                mapping[field] = idx
                break
    return mapping


def _safe_float(val: Any) -> Optional[float]:
    """Safely convert a value to float, stripping currency/percent symbols."""
    if val is None:
        return None
    if isinstance(val, (int, float)):
        return float(val)
    s = str(val).strip()
    if not s or s.lower() in ("n/a", "na", "-", "--", "null", "none", ""):
        return None
    s = re.sub(r"[$,\s%]", "", s)
    try:
        return float(s)
    except (ValueError, TypeError):
        return None


def _safe_str(val: Any) -> str:
    """Safely convert to string, handling None."""
    if val is None:
        return ""
    return str(val).strip()


# ═══════════════════════════════════════════════════════════════════════════════
# CHANNEL NAME -> PLATFORM KEY MAPPING
# ═══════════════════════════════════════════════════════════════════════════════

_CHANNEL_TO_PLATFORM: Dict[str, str] = {
    "google": "google_search",
    "google ads": "google_search",
    "google search": "google_search",
    "search": "google_search",
    "sem": "google_search",
    "ppc": "google_search",
    "google display": "programmatic",
    "gdn": "programmatic",
    "facebook": "meta_facebook",
    "meta": "meta_facebook",
    "meta ads": "meta_facebook",
    "instagram": "meta_instagram",
    "ig": "meta_instagram",
    "linkedin": "linkedin",
    "linkedin ads": "linkedin",
    "indeed": "indeed",
    "indeed sponsored": "indeed",
    "ziprecruiter": "indeed",
    "glassdoor": "indeed",
    "programmatic": "programmatic",
    "programmatic display": "programmatic",
    "job board": "indeed",
    "job boards": "indeed",
    "social": "meta_facebook",
    "social media": "meta_facebook",
    "display": "programmatic",
    "display ads": "programmatic",
    "career site": "google_search",
    "organic": "google_search",
    "appcast": "programmatic",
    "pandologic": "programmatic",
    "joveo": "programmatic",
    "recruitics": "programmatic",
    "talroo": "programmatic",
    "monster": "indeed",
    "careerbuilder": "indeed",
    "craigslist": "indeed",
    "snagajob": "indeed",
    "handshake": "indeed",
    "wellfound": "linkedin",
    "stackoverflow": "linkedin",
    "github jobs": "linkedin",
    "dice": "linkedin",
    "hired": "linkedin",
    "tiktok": "meta_facebook",
    "snapchat": "meta_facebook",
    "youtube": "programmatic",
    "spotify": "programmatic",
    "twitter": "meta_facebook",
    "x ads": "meta_facebook",
    "reddit": "meta_facebook",
    "niche board": "indeed",
    "specialty board": "indeed",
    "trade publication": "programmatic",
}


def _resolve_platform(channel_name: str) -> str:
    """Resolve a channel name to a trend_engine platform key."""
    ch = channel_name.lower().strip()
    for pattern, platform in _CHANNEL_TO_PLATFORM.items():
        if pattern in ch or ch in pattern:
            return platform
    return "programmatic"


# ═══════════════════════════════════════════════════════════════════════════════
# FALLBACK BENCHMARKS (used when trend_engine unavailable)
# ═══════════════════════════════════════════════════════════════════════════════

# Hardcoded fallback (kept for resilience if benchmark_registry unavailable)
_FALLBACK_BENCHMARKS: Dict[str, Dict[str, float]] = {
    "google_search": {"cpc": 2.69, "cpa": 45.00, "ctr": 0.042, "cpm": 10.00},
    "meta_facebook": {"cpc": 1.72, "cpa": 30.00, "ctr": 0.012, "cpm": 7.50},
    "meta_instagram": {"cpc": 1.50, "cpa": 35.00, "ctr": 0.010, "cpm": 8.00},
    "linkedin": {
        "cpc": 5.26,
        "cpa": 45.00,
        "ctr": 0.008,
        "cpm": 35.00,
    },  # Sponsored Jobs CPA $30-$90, US avg $45 (2026-04-07)
    "indeed": {"cpc": 0.50, "cpa": 25.00, "ctr": 0.040, "cpm": 5.00},
    "programmatic": {"cpc": 0.63, "cpa": 22.00, "ctr": 0.025, "cpm": 4.50},
}


def _get_fallback(platform: str, metric: str) -> float:
    """Get fallback benchmark value.

    Prefers benchmark_registry when available, else uses hardcoded fallback.
    """
    if _HAS_BENCHMARK_REGISTRY:
        return get_benchmark_value(platform, metric)
    return _FALLBACK_BENCHMARKS.get(platform, _FALLBACK_BENCHMARKS["programmatic"]).get(
        metric, 1.00
    )


# ═══════════════════════════════════════════════════════════════════════════════
# RECOMMENDED CHANNELS per collar type (for missing channel analysis)
# ═══════════════════════════════════════════════════════════════════════════════

_RECOMMENDED_CHANNELS: Dict[str, Dict[str, Dict[str, Any]]] = {
    "blue_collar": {
        "Indeed": {
            "platform": "indeed",
            "importance": "critical",
            "reason": "Indeed dominates blue-collar job search with 78% mobile apply rate. Cost-effective CPC model ideal for high-volume hourly roles.",
        },
        "Programmatic (Appcast/Joveo)": {
            "platform": "programmatic",
            "importance": "critical",
            "reason": "CPA-based programmatic campaigns auto-optimize bids across 100+ job sites, reducing cost-per-application by 30-50% vs manual posting.",
        },
        "Facebook/Meta": {
            "platform": "meta_facebook",
            "importance": "high",
            "reason": "Meta's broad reach and geo-targeting excels for local hourly hiring. 65% of blue-collar candidates are active Facebook users.",
        },
        "Google Ads": {
            "platform": "google_search",
            "importance": "medium",
            "reason": "Google search captures high-intent candidates actively searching for 'jobs near me'. Strong for warehouse, driver, and trade roles.",
        },
        "ZipRecruiter": {
            "platform": "indeed",
            "importance": "medium",
            "reason": "ZipRecruiter's AI matching and one-click apply drives strong application volume for entry-level and skilled trade positions.",
        },
        "Snagajob": {
            "platform": "indeed",
            "importance": "medium",
            "reason": "Specialized hourly job board with strong candidate pool for retail, food service, and warehouse roles.",
        },
    },
    "white_collar": {
        "LinkedIn": {
            "platform": "linkedin",
            "importance": "critical",
            "reason": "LinkedIn is the primary professional network with 950M+ members. InMail response rates are 3x higher than job board applications for white-collar roles.",
        },
        "Programmatic (Appcast/Joveo)": {
            "platform": "programmatic",
            "importance": "high",
            "reason": "Programmatic aggregation reaches passive candidates across niche professional job boards at optimized CPAs.",
        },
        "Google Ads": {
            "platform": "google_search",
            "importance": "high",
            "reason": "Google search captures professionals actively exploring new opportunities. Strong for tech, finance, and management roles.",
        },
        "Indeed": {
            "platform": "indeed",
            "importance": "high",
            "reason": "Indeed remains the #1 job site by volume. Sponsored listings ensure visibility for professional roles in competitive markets.",
        },
        "Glassdoor": {
            "platform": "indeed",
            "importance": "medium",
            "reason": "White-collar candidates research employer reviews before applying. Glassdoor presence improves application quality by 20-30%.",
        },
        "Employer Branding (Career Site/EVP)": {
            "platform": "google_search",
            "importance": "medium",
            "reason": "Investing in employer brand drives organic applicant flow and reduces dependency on paid channels by 15-25% over 6 months.",
        },
    },
    "grey_collar": {
        "Niche Job Boards": {
            "platform": "indeed",
            "importance": "critical",
            "reason": "Specialized boards (Vivian, NurseFly, HealthJobsNationwide) convert 2x better than general boards for healthcare and technical roles.",
        },
        "Programmatic (Appcast/Joveo)": {
            "platform": "programmatic",
            "importance": "critical",
            "reason": "Programmatic reach across specialty boards is essential for hard-to-fill grey-collar positions with critical talent shortages.",
        },
        "Indeed": {
            "platform": "indeed",
            "importance": "high",
            "reason": "Indeed's broad reach supplements niche boards for allied health, technician, and skilled service roles.",
        },
        "LinkedIn": {
            "platform": "linkedin",
            "importance": "medium",
            "reason": "LinkedIn targets credentialed professionals (RNs, therapists, certified technicians) who maintain professional profiles.",
        },
        "Facebook/Meta": {
            "platform": "meta_facebook",
            "importance": "medium",
            "reason": "Meta geo-targeting reaches local healthcare and technical workers. Effective for relocation campaigns and sign-on bonus promotion.",
        },
        "Google Ads": {
            "platform": "google_search",
            "importance": "medium",
            "reason": "Search ads capture nurses, technicians, and therapists searching for '[specialty] jobs near me' with high conversion intent.",
        },
    },
    "mixed": {
        "Programmatic (Appcast/Joveo)": {
            "platform": "programmatic",
            "importance": "critical",
            "reason": "Programmatic is the universal foundation for any multi-role campaign, automating bid optimization across the full job board ecosystem.",
        },
        "Indeed": {
            "platform": "indeed",
            "importance": "critical",
            "reason": "Indeed is the #1 job site globally. Essential for any recruitment media plan regardless of role type.",
        },
        "LinkedIn": {
            "platform": "linkedin",
            "importance": "high",
            "reason": "LinkedIn targets professional and semi-professional candidates with precision targeting by skills, experience, and industry.",
        },
        "Google Ads": {
            "platform": "google_search",
            "importance": "high",
            "reason": "Google captures high-intent job seekers across all collar types. Essential for branded and non-branded search campaigns.",
        },
        "Facebook/Meta": {
            "platform": "meta_facebook",
            "importance": "medium",
            "reason": "Meta's scale and targeting reach passive candidates. Particularly effective for local hiring and employer brand awareness.",
        },
        "Employer Branding": {
            "platform": "google_search",
            "importance": "medium",
            "reason": "Building employer brand reduces long-term cost-per-hire by 20-40% and improves quality of applicants.",
        },
    },
}


# ═══════════════════════════════════════════════════════════════════════════════
# 1. PARSE MEDIA PLAN
# ═══════════════════════════════════════════════════════════════════════════════


def parse_media_plan(file_bytes: bytes, filename: str) -> List[Dict[str, Any]]:
    """Parse uploaded Excel/CSV media plan into structured line items.

    Flexible column matching for: Channel/Platform, Budget/Spend,
    CPC, CPA, Target Roles, Target Locations, Duration.

    Returns list of dicts, one per channel/row with parsed values.
    Returns empty list on failure (never raises).
    """
    MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB
    if len(file_bytes) > MAX_FILE_SIZE:
        return []
    try:
        ext = os.path.splitext(os.path.basename(filename))[1].lstrip(".").lower()
        if ext in ("xlsx", "xls"):
            return _parse_excel(file_bytes)
        elif ext == "csv":
            return _parse_csv(file_bytes)
        else:
            logger.warning("Unsupported file type for media plan: %s", filename)
            return []
    except Exception as exc:
        logger.exception("Failed to parse media plan from %s: %s", filename, exc)
        return []


def _parse_excel(file_bytes: bytes) -> List[Dict[str, Any]]:
    """Parse Excel file into media plan line items."""
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)

    try:
        # Try all sheets, pick the one with the most parseable data
        best_records: List[Dict[str, Any]] = []

        for ws in wb.worksheets:
            rows = list(ws.iter_rows(values_only=True))
            if len(rows) < 2:
                continue

            # Try to find header row (first row with recognizable column names)
            for header_idx in range(min(5, len(rows))):
                candidate_headers = [
                    str(c).strip() if c else "" for c in rows[header_idx]
                ]
                col_map = _map_columns(candidate_headers)
                # Must have at least channel + budget
                if (
                    col_map.get("channel") is not None
                    and col_map.get("budget") is not None
                ):
                    records = _rows_to_line_items(
                        candidate_headers, rows[header_idx + 1 :]
                    )
                    if len(records) > len(best_records):
                        best_records = records
                    break

        return best_records
    finally:
        wb.close()


def _parse_csv(file_bytes: bytes) -> List[Dict[str, Any]]:
    """Parse CSV file into media plan line items."""
    import csv as csv_mod

    text = None
    for enc in ("utf-8", "utf-8-sig", "latin-1", "cp1252"):
        try:
            text = file_bytes.decode(enc)
            break
        except (UnicodeDecodeError, ValueError):
            continue
    if not text:
        return []
    reader = csv_mod.reader(io.StringIO(text))
    rows = list(reader)
    if len(rows) < 2:
        return []

    # Try first few rows as header
    for header_idx in range(min(5, len(rows))):
        candidate_headers = [str(c).strip() if c else "" for c in rows[header_idx]]
        col_map = _map_columns(candidate_headers)
        if col_map.get("channel") is not None and col_map.get("budget") is not None:
            return _rows_to_line_items(candidate_headers, rows[header_idx + 1 :])

    # Fallback: use first row
    headers = [str(c).strip() if c else "" for c in rows[0]]
    return _rows_to_line_items(headers, rows[1:])


def _rows_to_line_items(headers: List[str], data_rows: List) -> List[Dict[str, Any]]:
    """Convert parsed rows to list of media plan line item dicts."""
    col_map = _map_columns(headers)
    items: List[Dict[str, Any]] = []

    for row in data_rows:
        cells = list(row)
        # Skip empty rows
        if not any(c for c in cells if c is not None and str(c).strip()):
            continue

        ch_idx = col_map.get("channel")
        if ch_idx is None or ch_idx >= len(cells) or not cells[ch_idx]:
            continue

        channel = _safe_str(cells[ch_idx])
        if not channel or channel.lower() in (
            "total",
            "grand total",
            "sum",
            "subtotal",
        ):
            continue

        def _get_float(field: str) -> Optional[float]:
            idx = col_map.get(field)
            if idx is not None and idx < len(cells):
                return _safe_float(cells[idx])
            return None

        def _get_str(field: str) -> str:
            idx = col_map.get(field)
            if idx is not None and idx < len(cells):
                return _safe_str(cells[idx])
            return ""

        budget = _get_float("budget") or 0.0
        cpc = _get_float("cpc")
        cpa = _get_float("cpa")
        impressions = _get_float("impressions")
        clicks = _get_float("clicks")

        # Compute CPC from budget and clicks if not provided
        if cpc is None and clicks and clicks > 0 and budget > 0:
            cpc = round(budget / clicks, 2)

        # Compute CPA from budget and estimated applications if possible
        # (applications not usually in a plan, but clicks and apply rate might be derivable)

        item: Dict[str, Any] = {
            "channel": channel,
            "budget": budget,
            "cpc": cpc,
            "cpa": cpa,
            "impressions": impressions,
            "clicks": clicks,
            "target_roles": _get_str("target_roles"),
            "target_locations": _get_str("target_locations"),
            "duration": _get_str("duration"),
            "platform": _resolve_platform(channel),
        }
        items.append(item)

    return items


# ═══════════════════════════════════════════════════════════════════════════════
# 2. GET BENCHMARKS
# ═══════════════════════════════════════════════════════════════════════════════


def _get_benchmarks_for_platform(
    platform: str,
    industry: str = "general_entry_level",
    collar_type: str = "mixed",
    location: str = "",
) -> Dict[str, float]:
    """Pull benchmark CPC/CPA/CTR from trend_engine for a single platform."""
    bench: Dict[str, float] = {}
    now = datetime.datetime.now()
    month = now.month
    year = now.year

    for metric in ("cpc", "cpa", "ctr", "cpm"):
        if _HAS_TREND_ENGINE:
            try:
                result = _trend_engine.get_benchmark(
                    platform=platform,
                    industry=industry,
                    metric=metric,
                    collar_type=collar_type,
                    location=location,
                    month=month,
                    year=min(year, 2026),
                )
                bench[metric] = result.get("value", 0.0)
                bench[f"{metric}_confidence"] = result.get("data_confidence", 0.5)
            except Exception:
                bench[metric] = _get_fallback(platform, metric)
        else:
            bench[metric] = _get_fallback(platform, metric)

    return bench


def _get_all_benchmarks(
    industry: str = "general_entry_level",
    collar_type: str = "mixed",
    location: str = "",
) -> Dict[str, Dict[str, float]]:
    """Pull benchmarks for all platforms."""
    platforms = [
        "google_search",
        "meta_facebook",
        "linkedin",
        "indeed",
        "programmatic",
        "meta_instagram",
    ]
    return {
        p: _get_benchmarks_for_platform(p, industry, collar_type, location)
        for p in platforms
    }


# ═══════════════════════════════════════════════════════════════════════════════
# 3. AUDIT A SINGLE LINE ITEM
# ═══════════════════════════════════════════════════════════════════════════════


def _score_to_grade(score: float) -> str:
    """Convert efficiency score (0-100) to letter grade."""
    if score >= 85:
        return "A"
    elif score >= 70:
        return "B"
    elif score >= 55:
        return "C"
    elif score >= 40:
        return "D"
    else:
        return "F"


def audit_line_item(
    item: Dict[str, Any],
    industry: str = "general_entry_level",
    roles: Optional[List[str]] = None,
    locations: Optional[List[str]] = None,
    collar_type: str = "mixed",
    benchmarks_cache: Optional[Dict[str, Dict[str, float]]] = None,
) -> Dict[str, Any]:
    """Audit a single media plan line item against benchmarks.

    Returns dict with:
        benchmark_cpc, benchmark_cpa, variance_pct (CPC), cpa_variance_pct,
        grade (A-F), finding (overspend/underspend/optimal/no_data),
        savings_potential_dollars, efficiency_score (0-100)
    """
    platform = item.get("platform", _resolve_platform(item.get("channel") or ""))
    budget = item.get("budget", 0.0)
    planned_cpc = item.get("cpc")
    planned_cpa = item.get("cpa")

    # Get benchmark
    if benchmarks_cache and platform in benchmarks_cache:
        bench = benchmarks_cache[platform]
    else:
        location = ""
        if locations:
            location = (
                locations[0] if isinstance(locations[0], str) else str(locations[0])
            )
        bench = _get_benchmarks_for_platform(platform, industry, collar_type, location)

    benchmark_cpc = bench.get("cpc", 0.0)
    benchmark_cpa = bench.get("cpa", 0.0)
    benchmark_ctr = bench.get("ctr", 0.0)

    result: Dict[str, Any] = {
        "channel": item.get("channel", "Unknown"),
        "platform": platform,
        "planned_budget": budget,
        "planned_cpc": planned_cpc,
        "planned_cpa": planned_cpa,
        "benchmark_cpc": round(benchmark_cpc, 2),
        "benchmark_cpa": round(benchmark_cpa, 2),
        "benchmark_ctr": round(benchmark_ctr, 4),
        "cpc_variance_pct": 0.0,
        "cpa_variance_pct": 0.0,
        "finding": "no_data",
        "grade": "C",
        "efficiency_score": 50.0,
        "savings_potential": 0.0,
        "detail": "",
    }

    scores: List[float] = []

    # CPC analysis
    if planned_cpc is not None and benchmark_cpc > 0:
        cpc_var = ((planned_cpc - benchmark_cpc) / benchmark_cpc) * 100
        result["cpc_variance_pct"] = round(cpc_var, 1)
        # Lower CPC is better: negative variance = good
        cpc_score = min(100, max(0, 50 - cpc_var * 2))
        scores.append(cpc_score)

        if cpc_var > 30:
            result[
                "detail"
            ] += f"CPC is {cpc_var:.0f}% above benchmark (${planned_cpc:.2f} vs ${benchmark_cpc:.2f}). "
        elif cpc_var < -20:
            result[
                "detail"
            ] += f"CPC is {abs(cpc_var):.0f}% below benchmark -- strong value. "
    else:
        scores.append(50.0)

    # CPA analysis
    if planned_cpa is not None and benchmark_cpa > 0:
        cpa_var = ((planned_cpa - benchmark_cpa) / benchmark_cpa) * 100
        result["cpa_variance_pct"] = round(cpa_var, 1)
        cpa_score = min(100, max(0, 50 - cpa_var * 2))
        scores.append(cpa_score)

        if cpa_var > 30:
            result[
                "detail"
            ] += f"CPA is {cpa_var:.0f}% above benchmark (${planned_cpa:.2f} vs ${benchmark_cpa:.2f}). "
        elif cpa_var < -20:
            result[
                "detail"
            ] += f"CPA is {abs(cpa_var):.0f}% below benchmark -- efficient targeting. "
    else:
        scores.append(50.0)

    # Budget proportionality score (penalize extremely high concentrations)
    if budget > 0 and planned_cpc is not None and benchmark_cpc > 0:
        budget_score = min(
            100,
            max(
                0, 60 - abs(((planned_cpc - benchmark_cpc) / benchmark_cpc) * 100) * 1.5
            ),
        )
        scores.append(budget_score)
    else:
        scores.append(50.0)

    # Overall efficiency score (CPC 40%, CPA 40%, Budget 20%)
    if len(scores) >= 3:
        efficiency = scores[0] * 0.40 + scores[1] * 0.40 + scores[2] * 0.20
    elif scores:
        efficiency = sum(scores) / len(scores)
    else:
        efficiency = 50.0

    result["efficiency_score"] = round(min(100, max(0, efficiency)), 1)
    result["grade"] = _score_to_grade(result["efficiency_score"])

    # Determine finding
    if planned_cpc is None and planned_cpa is None:
        result["finding"] = "no_data"
        result["detail"] = (
            "No CPC or CPA data provided. Cannot fully assess cost efficiency."
        )
    elif result["efficiency_score"] >= 70:
        result["finding"] = "optimal"
        if not result["detail"]:
            result["detail"] = (
                "Planned costs are at or below benchmark levels. Good allocation."
            )
    elif result["efficiency_score"] >= 40:
        result["finding"] = "review"
        if not result["detail"]:
            result["detail"] = (
                "Costs are slightly above benchmarks. Consider negotiating rates."
            )
    else:
        result["finding"] = "overspend"
        if not result["detail"]:
            result["detail"] = (
                "Significant overspend detected. Costs are well above industry benchmarks."
            )

    # Calculate savings potential
    if (
        planned_cpc is not None
        and benchmark_cpc > 0
        and planned_cpc > benchmark_cpc
        and budget > 0
    ):
        # How much could be saved if CPC were at benchmark
        estimated_clicks = budget / planned_cpc if planned_cpc > 0 else 0
        cost_at_benchmark = estimated_clicks * benchmark_cpc
        result["savings_potential"] = round(max(0, budget - cost_at_benchmark), 2)
    elif (
        planned_cpa is not None
        and benchmark_cpa > 0
        and planned_cpa > benchmark_cpa
        and budget > 0
    ):
        estimated_apps = budget / planned_cpa if planned_cpa > 0 else 0
        cost_at_benchmark = estimated_apps * benchmark_cpa
        result["savings_potential"] = round(max(0, budget - cost_at_benchmark), 2)

    return result


# ═══════════════════════════════════════════════════════════════════════════════
# 4. IDENTIFY MISSING CHANNELS
# ═══════════════════════════════════════════════════════════════════════════════


def identify_missing_channels(
    current_channels: List[str],
    industry: str = "general_entry_level",
    roles: Optional[List[str]] = None,
    collar_type: str = "mixed",
) -> List[Dict[str, Any]]:
    """Compare client's channel list against recommended channels.

    Flags channels they SHOULD be using but aren't. Uses collar_intelligence
    strategy data when available, falls back to built-in recommendations.

    Returns list of dicts with: channel, importance, reason, platform.
    """
    # Determine collar type from roles if not provided
    if collar_type == "mixed" and roles and _HAS_COLLAR_INTEL:
        try:
            result = _collar_intel.classify_collar(roles[0], industry=industry)
            collar_type = result.get("collar_type", "mixed")
        except Exception:
            pass

    # Normalize current channel names for matching
    current_normalized = set()
    for ch in current_channels:
        ch_lower = ch.lower().strip()
        current_normalized.add(ch_lower)
        # Also add the platform key for broader matching
        current_normalized.add(_resolve_platform(ch))

    # Get recommended channels for this collar type
    recommended = _RECOMMENDED_CHANNELS.get(collar_type, _RECOMMENDED_CHANNELS["mixed"])

    missing: List[Dict[str, Any]] = []
    for channel_name, info in recommended.items():
        # Check if this channel (or its platform) is already in the plan
        ch_lower = channel_name.lower().strip()
        platform = info.get("platform") or ""

        # Flexible matching: check if any current channel maps to same platform
        # or has similar name
        is_present = False
        for curr in current_normalized:
            if ch_lower in curr or curr in ch_lower:
                is_present = True
                break
            if platform and platform == curr:
                is_present = True
                break

        if not is_present:
            missing.append(
                {
                    "channel": channel_name,
                    "platform": platform,
                    "importance": info.get("importance", "medium"),
                    "reason": info.get(
                        "reason", "Recommended for your recruitment strategy."
                    ),
                }
            )

    # Sort by importance
    importance_order = {"critical": 0, "high": 1, "medium": 2, "low": 3}
    missing.sort(key=lambda x: importance_order.get(x["importance"], 3))

    return missing


# ═══════════════════════════════════════════════════════════════════════════════
# 5. CALCULATE TOTAL SAVINGS POTENTIAL
# ═══════════════════════════════════════════════════════════════════════════════


def calculate_savings_potential(audit_results: List[Dict[str, Any]]) -> Dict[str, Any]:
    """Total potential savings across all overspend channels.

    Returns summary dict with total savings, per-channel breakdown,
    and narrative text.
    """
    total_savings = 0.0
    total_budget = 0.0
    channel_savings: List[Dict[str, Any]] = []
    overspend_count = 0

    for ar in audit_results:
        budget = ar.get("planned_budget", 0.0)
        savings = ar.get("savings_potential", 0.0)
        total_budget += budget
        total_savings += savings

        if savings > 0:
            overspend_count += 1
            channel_savings.append(
                {
                    "channel": ar.get("channel", "Unknown"),
                    "planned_budget": budget,
                    "savings": savings,
                    "savings_pct": round(
                        (savings / budget * 100) if budget > 0 else 0, 1
                    ),
                }
            )

    channel_savings.sort(key=lambda x: x["savings"], reverse=True)

    savings_pct_of_total = round(
        (total_savings / total_budget * 100) if total_budget > 0 else 0, 1
    )

    if total_savings >= 10000:
        narrative = (
            f"Your media plan has ${total_savings:,.0f} in potential savings ({savings_pct_of_total}% of total budget). "
            f"{overspend_count} channel(s) are priced above industry benchmarks. "
            f"Negotiate vendor rates or reallocate to more cost-effective platforms."
        )
    elif total_savings > 0:
        narrative = (
            f"Potential savings of ${total_savings:,.0f} identified ({savings_pct_of_total}% of budget). "
            f"Minor optimizations in {overspend_count} channel(s) could improve cost efficiency."
        )
    else:
        narrative = (
            "Your planned costs are in line with or below industry benchmarks. "
            "Focus on channel coverage and allocation balance for maximum impact."
        )

    return {
        "total_savings": round(total_savings, 2),
        "total_budget": round(total_budget, 2),
        "savings_pct": savings_pct_of_total,
        "overspend_channels": overspend_count,
        "channel_savings": channel_savings[:10],
        "narrative": narrative,
    }


# ═══════════════════════════════════════════════════════════════════════════════
# 6. GENERATE RECOMMENDATIONS
# ═══════════════════════════════════════════════════════════════════════════════


def generate_recommendations(
    audit_results: List[Dict[str, Any]],
    missing_channels: List[Dict[str, Any]],
    industry: str = "general_entry_level",
) -> List[Dict[str, Any]]:
    """Prioritized list of actions: cut overspend, add missing channels, rebalance.

    Returns list of recommendation dicts, each with:
        priority (high/medium/low), action, channel, detail, category
    """
    recommendations: List[Dict[str, Any]] = []
    priority_counter = 1

    # 1. Flag overspend channels (highest priority)
    overspend_items = [ar for ar in audit_results if ar.get("finding") == "overspend"]
    overspend_items.sort(key=lambda x: x.get("savings_potential") or 0, reverse=True)

    for ar in overspend_items:
        savings = ar.get("savings_potential") or 0
        recommendations.append(
            {
                "priority": "high",
                "priority_num": priority_counter,
                "action": f"Reduce {ar['channel']} CPC/CPA to benchmark levels",
                "channel": ar.get("channel") or "",
                "detail": (
                    f"Planned CPC ${ar.get('planned_cpc') or 0:.2f} is {ar.get('cpc_variance_pct') or 0:.0f}% above "
                    f"benchmark ${ar.get('benchmark_cpc') or 0:.2f}. "
                    f"Potential savings: ${savings:,.0f}. "
                    f"Negotiate rates, optimize targeting, or consider alternative vendors."
                ),
                "category": "overspend",
                "savings": savings,
            }
        )
        priority_counter += 1

    # 2. Add missing critical/high channels
    for mc in missing_channels:
        if mc["importance"] in ("critical", "high"):
            recommendations.append(
                {
                    "priority": "high" if mc["importance"] == "critical" else "medium",
                    "priority_num": priority_counter,
                    "action": f"Add {mc['channel']} to your media plan",
                    "channel": mc["channel"],
                    "detail": mc["reason"],
                    "category": "missing_channel",
                    "savings": 0,
                }
            )
            priority_counter += 1

    # 3. Channels needing review
    review_items = [ar for ar in audit_results if ar.get("finding") == "review"]
    for ar in review_items:
        recommendations.append(
            {
                "priority": "medium",
                "priority_num": priority_counter,
                "action": f"Review {ar['channel']} pricing",
                "channel": ar.get("channel") or "",
                "detail": (
                    f"{ar.get('channel') or ''} costs are slightly above benchmarks (Grade: {ar.get('grade', 'C')}). "
                    f"{ar.get('detail') or ''}"
                ),
                "category": "review",
                "savings": ar.get("savings_potential") or 0,
            }
        )
        priority_counter += 1

    # 4. Missing medium-importance channels
    for mc in missing_channels:
        if mc["importance"] == "medium":
            recommendations.append(
                {
                    "priority": "low",
                    "priority_num": priority_counter,
                    "action": f"Consider adding {mc['channel']}",
                    "channel": mc["channel"],
                    "detail": mc["reason"],
                    "category": "missing_channel",
                    "savings": 0,
                }
            )
            priority_counter += 1

    # 5. Budget rebalancing suggestion (if we have multiple channels)
    optimal_items = [ar for ar in audit_results if ar.get("finding") == "optimal"]
    if optimal_items and overspend_items:
        best_channel = max(optimal_items, key=lambda x: x.get("efficiency_score") or 0)
        recommendations.append(
            {
                "priority": "medium",
                "priority_num": priority_counter,
                "action": f"Reallocate overspend budget to {best_channel['channel']}",
                "channel": best_channel.get("channel") or "",
                "detail": (
                    f"{best_channel['channel']} has the best efficiency score ({best_channel.get('efficiency_score') or 0:.0f}/100, "
                    f"Grade {best_channel.get('grade', 'C')}). "
                    f"Move savings from overpriced channels here for maximum ROI."
                ),
                "category": "rebalance",
                "savings": 0,
            }
        )
        priority_counter += 1

    # 6. Industry-specific advice
    industry_tips = _get_industry_tips(industry)
    if industry_tips:
        recommendations.append(
            {
                "priority": "low",
                "priority_num": priority_counter,
                "action": "Industry best practice",
                "channel": "",
                "detail": industry_tips,
                "category": "industry_insight",
                "savings": 0,
            }
        )

    return recommendations[:15]  # Cap at 15


def _get_industry_tips(industry: str) -> str:
    """Return industry-specific recruitment advertising tips."""
    tips = {
        "healthcare_medical": "Healthcare: Prioritize niche boards (Vivian, NurseFly) and sign-on bonus campaigns. Travel nursing demand is seasonal -- Q1 and Q4 see highest demand.",
        "blue_collar_trades": "Blue Collar: Mobile-first apply experience is critical (78% apply via mobile). Indeed and programmatic dominate. Keep applications under 2 minutes.",
        "tech_engineering": "Tech: Developer communities (GitHub, Stack Overflow) outperform LinkedIn for IC roles. Technical assessment platforms improve quality-of-hire metrics.",
        "finance_banking": "Finance: LinkedIn InMail and niche boards (eFinancialCareers) drive highest quality. Compliance messaging is essential in regulated markets.",
        "logistics_supply_chain": "Logistics: High-volume programmatic with CPA caps is most efficient. Geo-fence distribution centers and competitor locations for targeted reach.",
        "retail_consumer": "Retail: Seasonal hiring requires 60-90 day lead time. Facebook and Indeed dominate hourly retail recruitment. Volume discounts on job board packages reduce CPA 20-30%.",
        "hospitality_travel": "Hospitality: Mobile apply and social media campaigns drive best results for hourly roles. Partner with hospitality-specific boards for management positions.",
        "aerospace_defense": "Aerospace: Specialized engineering boards and trade show presence are critical. Security clearance roles require targeted sourcing on ClearanceJobs and similar platforms.",
        "construction_real_estate": "Construction: Local job boards and trade union partnerships outperform national platforms. Safety certification targeting improves applicant quality.",
    }
    return tips.get(industry, "")


# ═══════════════════════════════════════════════════════════════════════════════
# 7. GENERATE AUDIT SCORECARD
# ═══════════════════════════════════════════════════════════════════════════════


def generate_audit_scorecard(
    audit_results: List[Dict[str, Any]],
    missing_channels: List[Dict[str, Any]],
) -> Dict[str, Any]:
    """Overall audit grade (A-F), budget efficiency score (0-100), channel coverage score.

    Returns comprehensive scorecard dict.
    """
    if not audit_results:
        return {
            "overall_grade": "N/A",
            "budget_efficiency_score": 0,
            "channel_coverage_score": 0,
            "channel_count": 0,
            "channel_grades": [],
            "findings_summary": {},
        }

    # Budget efficiency score: weighted average of line item scores
    total_budget = sum(ar.get("planned_budget") or 0 for ar in audit_results)
    if total_budget > 0:
        weighted_score = (
            sum(
                ar.get("efficiency_score", 50) * ar.get("planned_budget") or 0
                for ar in audit_results
            )
            / total_budget
        )
    else:
        weighted_score = sum(
            ar.get("efficiency_score", 50) for ar in audit_results
        ) / len(audit_results)

    budget_efficiency_score = round(min(100, max(0, weighted_score)), 1)

    # Channel coverage score: based on how many recommended channels are present
    total_recommended = len(missing_channels) + len(audit_results)
    if total_recommended > 0:
        coverage_ratio = len(audit_results) / total_recommended
        channel_coverage_score = round(min(100, coverage_ratio * 100), 1)
    else:
        channel_coverage_score = 50.0

    # Boost coverage score if critical channels are not missing
    critical_missing = sum(
        1 for mc in missing_channels if mc.get("importance") == "critical"
    )
    if critical_missing > 0:
        channel_coverage_score = max(0, channel_coverage_score - critical_missing * 15)

    # Overall grade: 60% efficiency + 40% coverage
    overall_score = budget_efficiency_score * 0.60 + channel_coverage_score * 0.40
    overall_grade = _score_to_grade(overall_score)

    # Findings summary
    findings_count: Dict[str, int] = {
        "optimal": 0,
        "review": 0,
        "overspend": 0,
        "no_data": 0,
    }
    for ar in audit_results:
        finding = ar.get("finding", "no_data")
        findings_count[finding] = findings_count.get(finding, 0) + 1

    # Channel grades
    channel_grades = [
        {
            "channel": ar.get("channel") or "",
            "grade": ar.get("grade", "C"),
            "score": ar.get("efficiency_score", 50),
            "finding": ar.get("finding", "no_data"),
        }
        for ar in audit_results
    ]
    channel_grades.sort(key=lambda x: x["score"], reverse=True)

    return {
        "overall_grade": overall_grade,
        "overall_score": round(overall_score, 1),
        "budget_efficiency_score": budget_efficiency_score,
        "channel_coverage_score": round(channel_coverage_score, 1),
        "channel_count": len(audit_results),
        "missing_channel_count": len(missing_channels),
        "channel_grades": channel_grades,
        "findings_summary": findings_count,
    }


def _finding_text(finding: str) -> str:
    """Convert finding key to human-readable label."""
    return {
        "optimal": "Optimal",
        "review": "Needs Review",
        "overspend": "Overspend",
        "no_data": "Insufficient Data",
    }.get(finding, finding.replace("_", " ").title())


# ═══════════════════════════════════════════════════════════════════════════════
# 8. GENERATE AUDIT EXCEL
# ═══════════════════════════════════════════════════════════════════════════════


def generate_audit_excel(
    report_data: Dict[str, Any], client_name: str = "Client"
) -> bytes:
    """Generate 4-sheet Excel audit report.

    Sheets:
        1. Audit Summary
        2. Line-by-Line Analysis
        3. Missing Channels
        4. Recommendations

    Sapphire Blue palette, Calibri font, data starts at column B.
    """
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    except ImportError:
        logger.error("openpyxl not available for Excel generation")
        return b""

    # Design tokens -- Sapphire Blue palette
    NAVY = "0F172A"
    SAPPHIRE = "2563EB"
    BLUE_LIGHT = "DBEAFE"
    BLUE_PALE = "EFF6FF"
    STONE = "1C1917"
    MUTED = "78716C"
    WARM_GRAY = "E7E5E4"
    WHITE = "FFFFFF"
    GREEN = "16A34A"
    GREEN_BG = "DCFCE7"
    AMBER = "D97706"
    AMBER_BG = "FEF3C7"
    RED = "DC2626"
    RED_BG = "FEE2E2"

    # Fonts
    f_title = Font(name="Calibri", bold=True, size=18, color=NAVY)
    f_section = Font(name="Calibri", bold=True, size=14, color=WHITE)
    f_subsection = Font(name="Calibri", bold=True, size=12, color=NAVY)
    f_header = Font(name="Calibri", bold=True, size=10, color=WHITE)
    f_body = Font(name="Calibri", size=10, color=STONE)
    f_body_bold = Font(name="Calibri", bold=True, size=10, color=STONE)
    f_hero = Font(name="Calibri", bold=True, size=22, color=SAPPHIRE)
    f_hero_label = Font(name="Calibri", size=9, color=MUTED)
    f_metric_value = Font(name="Calibri", bold=True, size=14, color=NAVY)
    f_grade_big = Font(name="Calibri", bold=True, size=36, color=WHITE)
    f_footnote = Font(name="Calibri", italic=True, size=9, color=MUTED)
    f_green = Font(name="Calibri", bold=True, size=10, color=GREEN)
    f_red = Font(name="Calibri", bold=True, size=10, color=RED)
    f_amber = Font(name="Calibri", bold=True, size=10, color=AMBER)

    # Fills
    fill_navy = PatternFill(start_color=NAVY, end_color=NAVY, fill_type="solid")
    fill_sapphire = PatternFill(
        start_color=SAPPHIRE, end_color=SAPPHIRE, fill_type="solid"
    )
    fill_light = PatternFill(
        start_color=BLUE_LIGHT, end_color=BLUE_LIGHT, fill_type="solid"
    )
    fill_pale = PatternFill(
        start_color=BLUE_PALE, end_color=BLUE_PALE, fill_type="solid"
    )
    fill_white = PatternFill(start_color=WHITE, end_color=WHITE, fill_type="solid")
    fill_green_bg = PatternFill(
        start_color=GREEN_BG, end_color=GREEN_BG, fill_type="solid"
    )
    fill_amber_bg = PatternFill(
        start_color=AMBER_BG, end_color=AMBER_BG, fill_type="solid"
    )
    fill_red_bg = PatternFill(start_color=RED_BG, end_color=RED_BG, fill_type="solid")
    fill_green = PatternFill(start_color=GREEN, end_color=GREEN, fill_type="solid")
    fill_amber = PatternFill(start_color=AMBER, end_color=AMBER, fill_type="solid")
    fill_red = PatternFill(start_color=RED, end_color=RED, fill_type="solid")

    def _grade_fill(grade: str):
        if grade == "A":
            return fill_green_bg
        elif grade == "B":
            return PatternFill(
                start_color="DCFCE7", end_color="DCFCE7", fill_type="solid"
            )
        elif grade == "C":
            return fill_amber_bg
        elif grade == "D":
            return PatternFill(
                start_color="FED7AA", end_color="FED7AA", fill_type="solid"
            )
        else:
            return fill_red_bg

    def _grade_font(grade: str):
        if grade in ("A", "B"):
            return f_green
        elif grade == "C":
            return f_amber
        else:
            return f_red

    pass  # _finding_text is module-level

    # Alignment
    al_center = Alignment(horizontal="center", vertical="center", wrap_text=True)
    al_left = Alignment(horizontal="left", vertical="top", wrap_text=True)
    al_right = Alignment(horizontal="right", vertical="top", wrap_text=True)

    # Border
    border = Border(
        left=Side(style="thin", color=WARM_GRAY),
        right=Side(style="thin", color=WARM_GRAY),
        top=Side(style="thin", color=WARM_GRAY),
        bottom=Side(style="thin", color=WARM_GRAY),
    )

    COL_START = 2  # Column B

    wb = Workbook()

    scorecard = report_data.get("scorecard", {})
    audit_results = report_data.get("audit_results") or []
    missing_channels = report_data.get("missing_channels") or []
    recommendations = report_data.get("recommendations") or []
    savings = report_data.get("savings", {})

    # ── Helpers ─────────────────────────────────────────────────────────
    def _write_section_header(ws, row, title, col_start=COL_START, col_end=8):
        for c in range(col_start, col_end + 1):
            cell = ws.cell(row=row, column=c)
            cell.fill = fill_navy
        ws.cell(row=row, column=col_start, value=title).font = f_section
        ws.cell(row=row, column=col_start).alignment = al_left
        return row + 1

    def _write_table_header(ws, row, headers, col_start=COL_START):
        for i, h in enumerate(headers):
            cell = ws.cell(row=row, column=col_start + i, value=h)
            cell.font = f_header
            cell.fill = fill_sapphire
            cell.alignment = al_center
            cell.border = border
        return row + 1

    def _write_table_row(ws, row, values, col_start=COL_START, fonts=None, fills=None):
        for i, v in enumerate(values):
            cell = ws.cell(row=row, column=col_start + i, value=v)
            cell.font = fonts[i] if fonts and i < len(fonts) else f_body
            if fills and i < len(fills) and fills[i]:
                cell.fill = fills[i]
            cell.alignment = al_center if i > 0 else al_left
            cell.border = border
        return row + 1

    # ══════════════════════════════════════════════════════════════════
    # SHEET 1: Audit Summary
    # ══════════════════════════════════════════════════════════════════
    ws1 = wb.active
    ws1.title = "Audit Summary"
    ws1.sheet_properties.tabColor = NAVY

    ws1.column_dimensions["A"].width = 3
    for col_letter in ["B", "C", "D", "E", "F", "G", "H"]:
        ws1.column_dimensions[col_letter].width = 18

    row = 2
    ws1.cell(row=row, column=COL_START, value="Media Plan Audit Report").font = f_title
    row += 1
    ws1.cell(
        row=row,
        column=COL_START,
        value=f"{client_name} | Generated {datetime.datetime.now().strftime('%B %d, %Y')}",
    ).font = f_footnote
    row += 2

    # Overall Grade Card
    row = _write_section_header(ws1, row, "AUDIT SCORECARD")
    row += 1

    # Grade
    ws1.cell(row=row, column=COL_START, value="Overall Grade").font = f_subsection
    grade_cell = ws1.cell(
        row=row, column=COL_START + 1, value=scorecard.get("overall_grade", "N/A")
    )
    grade_cell.font = Font(name="Calibri", bold=True, size=28, color=WHITE)
    _g = scorecard.get("overall_grade", "C")
    if _g in ("A", "B"):
        grade_cell.fill = fill_green
    elif _g == "C":
        grade_cell.fill = fill_amber
    else:
        grade_cell.fill = fill_red
    grade_cell.alignment = al_center

    ws1.cell(row=row, column=COL_START + 2, value="Budget Efficiency").font = (
        f_hero_label
    )
    ws1.cell(
        row=row,
        column=COL_START + 3,
        value=f"{scorecard.get('budget_efficiency_score') or 0}/100",
    ).font = f_metric_value
    ws1.cell(row=row, column=COL_START + 4, value="Channel Coverage").font = (
        f_hero_label
    )
    ws1.cell(
        row=row,
        column=COL_START + 5,
        value=f"{scorecard.get('channel_coverage_score') or 0}/100",
    ).font = f_metric_value
    row += 2

    # Savings Banner
    total_savings = savings.get("total_savings") or 0
    total_budget = savings.get("total_budget") or 0
    row = _write_section_header(ws1, row, "SAVINGS OPPORTUNITY")
    row += 1
    if total_savings > 0:
        ws1.cell(
            row=row, column=COL_START, value=f"Potential Savings: ${total_savings:,.0f}"
        ).font = Font(name="Calibri", bold=True, size=16, color=GREEN)
        ws1.cell(
            row=row + 1,
            column=COL_START,
            value=f"{savings.get('savings_pct') or 0}% of total planned budget (${total_budget:,.0f})",
        ).font = f_footnote
    else:
        ws1.cell(
            row=row, column=COL_START, value="No significant overspend detected"
        ).font = Font(name="Calibri", bold=True, size=14, color=GREEN)
        ws1.cell(
            row=row + 1,
            column=COL_START,
            value="Your planned costs are in line with industry benchmarks.",
        ).font = f_footnote
    row += 3

    # Findings Summary
    findings = scorecard.get("findings_summary", {})
    row = _write_section_header(ws1, row, "FINDINGS SUMMARY")
    labels = ["Optimal", "Needs Review", "Overspend", "Insufficient Data"]
    values_list = [
        findings.get("optimal") or 0,
        findings.get("review") or 0,
        findings.get("overspend") or 0,
        findings.get("no_data") or 0,
    ]
    for i, (label, val) in enumerate(zip(labels, values_list)):
        col = COL_START + i
        ws1.cell(row=row, column=col, value=label).font = f_hero_label
        ws1.cell(row=row, column=col).alignment = al_center
        ws1.cell(row=row + 1, column=col, value=str(val)).font = f_metric_value
        ws1.cell(row=row + 1, column=col).alignment = al_center
        if i == 0:
            ws1.cell(row=row + 1, column=col).fill = fill_green_bg
        elif i == 1:
            ws1.cell(row=row + 1, column=col).fill = fill_amber_bg
        elif i == 2:
            ws1.cell(row=row + 1, column=col).fill = fill_red_bg
        else:
            ws1.cell(row=row + 1, column=col).fill = fill_light
        ws1.cell(row=row + 1, column=col).border = border
    row += 3

    # Channel Grades
    row = _write_section_header(ws1, row, "CHANNEL GRADES")
    row = _write_table_header(
        ws1, row, ["Channel", "Grade", "Score", "Finding", "CPC Variance", "Savings"]
    )
    for ar in audit_results:
        grade = ar.get("grade", "C")
        cpc_var = ar.get("cpc_variance_pct") or 0
        sav = ar.get("savings_potential") or 0

        def _var_str(v):
            arrow = "+" if v > 0 else ""
            return f"{arrow}{v:.1f}%"

        fonts_row = [
            f_body_bold,
            _grade_font(grade),
            f_body,
            f_body,
            f_green if cpc_var < 0 else f_red if cpc_var > 20 else f_amber,
            f_green if sav == 0 else f_red,
        ]
        fills_row = [None, _grade_fill(grade), None, None, None, None]
        row = _write_table_row(
            ws1,
            row,
            [
                ar.get("channel") or "",
                grade,
                f"{ar.get('efficiency_score') or 0:.0f}",
                _finding_text(ar.get("finding") or ""),
                _var_str(cpc_var),
                f"${sav:,.0f}" if sav > 0 else "-",
            ],
            fonts=fonts_row,
            fills=fills_row,
        )

    row += 1
    ws1.cell(
        row=row,
        column=COL_START,
        value="Negative CPC variance = below benchmark (good). Positive = above benchmark (overspend).",
    ).font = f_footnote

    # ══════════════════════════════════════════════════════════════════
    # SHEET 2: Line-by-Line Analysis
    # ══════════════════════════════════════════════════════════════════
    ws2 = wb.create_sheet("Line-by-Line Analysis")
    ws2.sheet_properties.tabColor = SAPPHIRE

    ws2.column_dimensions["A"].width = 3
    for col_letter in ["B", "C", "D", "E", "F", "G", "H", "I", "J"]:
        ws2.column_dimensions[col_letter].width = 16

    row = 2
    row = _write_section_header(ws2, row, "DETAILED LINE ITEM ANALYSIS", col_end=10)
    row = _write_table_header(
        ws2,
        row,
        [
            "Channel",
            "Planned Budget",
            "Planned CPC",
            "Benchmark CPC",
            "CPC Variance",
            "Planned CPA",
            "Benchmark CPA",
            "Grade",
            "Finding",
        ],
    )

    for ar in audit_results:
        grade = ar.get("grade", "C")
        cpc_var = ar.get("cpc_variance_pct") or 0
        planned_cpc = ar.get("planned_cpc")
        planned_cpa = ar.get("planned_cpa")

        fonts_row = [
            f_body_bold,
            f_body,
            f_body,
            f_body,
            f_green if cpc_var < 0 else f_red if cpc_var > 20 else f_amber,
            f_body,
            f_body,
            _grade_font(grade),
            f_body,
        ]
        fills_row = [None] * 7 + [_grade_fill(grade), None]

        row = _write_table_row(
            ws2,
            row,
            [
                ar.get("channel") or "",
                f"${ar.get('planned_budget') or 0:,.2f}",
                f"${planned_cpc:.2f}" if planned_cpc is not None else "N/A",
                f"${ar.get('benchmark_cpc') or 0:.2f}",
                f"{'+' if cpc_var > 0 else ''}{cpc_var:.1f}%",
                f"${planned_cpa:.2f}" if planned_cpa is not None else "N/A",
                f"${ar.get('benchmark_cpa') or 0:.2f}",
                grade,
                _finding_text(ar.get("finding") or ""),
            ],
            fonts=fonts_row,
            fills=fills_row,
        )

    row += 2
    # Detail notes
    row = _write_section_header(ws2, row, "ANALYSIS NOTES", col_end=10)
    for ar in audit_results:
        detail = ar.get("detail") or ""
        if detail:
            ws2.cell(row=row, column=COL_START, value=ar.get("channel") or "").font = (
                f_body_bold
            )
            ws2.cell(row=row, column=COL_START + 1, value=detail).font = f_body
            ws2.cell(row=row, column=COL_START + 1).alignment = al_left
            ws2.merge_cells(
                start_row=row,
                start_column=COL_START + 1,
                end_row=row,
                end_column=COL_START + 8,
            )
            row += 1

    # ══════════════════════════════════════════════════════════════════
    # SHEET 3: Missing Channels
    # ══════════════════════════════════════════════════════════════════
    ws3 = wb.create_sheet("Missing Channels")
    ws3.sheet_properties.tabColor = "D97706"

    ws3.column_dimensions["A"].width = 3
    ws3.column_dimensions["B"].width = 24
    ws3.column_dimensions["C"].width = 14
    ws3.column_dimensions["D"].width = 60

    row = 2
    row = _write_section_header(
        ws3, row, "RECOMMENDED CHANNELS NOT IN YOUR PLAN", col_end=4
    )
    row += 1

    if missing_channels:
        ws3.cell(
            row=row,
            column=COL_START,
            value=f"{len(missing_channels)} recommended channel(s) are missing from your current plan.",
        ).font = f_subsection
        row += 2

        row = _write_table_header(
            ws3, row, ["Channel", "Importance", "Why You Should Add This"]
        )

        for mc in missing_channels:
            importance = mc.get("importance", "medium")
            if importance == "critical":
                imp_font = f_red
                imp_fill = fill_red_bg
            elif importance == "high":
                imp_font = f_amber
                imp_fill = fill_amber_bg
            else:
                imp_font = f_body
                imp_fill = fill_light

            row = _write_table_row(
                ws3,
                row,
                [
                    mc.get("channel") or "",
                    importance.upper(),
                    mc.get("reason") or "",
                ],
                fonts=[f_body_bold, imp_font, f_body],
                fills=[None, imp_fill, None],
            )
    else:
        ws3.cell(
            row=row,
            column=COL_START,
            value="Your media plan covers all recommended channels. Excellent coverage!",
        ).font = Font(name="Calibri", bold=True, size=12, color=GREEN)

    # ══════════════════════════════════════════════════════════════════
    # SHEET 4: Recommendations
    # ══════════════════════════════════════════════════════════════════
    ws4 = wb.create_sheet("Recommendations")
    ws4.sheet_properties.tabColor = GREEN

    ws4.column_dimensions["A"].width = 3
    ws4.column_dimensions["B"].width = 6
    ws4.column_dimensions["C"].width = 12
    ws4.column_dimensions["D"].width = 35
    ws4.column_dimensions["E"].width = 55
    ws4.column_dimensions["F"].width = 14

    row = 2
    row = _write_section_header(ws4, row, "PRIORITIZED RECOMMENDATIONS", col_end=6)
    row += 1
    ws4.cell(
        row=row,
        column=COL_START,
        value="Actions ranked by impact. Start with high-priority items for maximum savings.",
    ).font = f_footnote
    row += 2

    row = _write_table_header(
        ws4, row, ["#", "Priority", "Action", "Detail", "Est. Savings"]
    )

    for rec in recommendations:
        priority = rec.get("priority", "medium")
        if priority == "high":
            pri_font = f_red
            pri_fill = fill_red_bg
        elif priority == "medium":
            pri_font = f_amber
            pri_fill = fill_amber_bg
        else:
            pri_font = f_body
            pri_fill = fill_light

        rec_savings = rec.get("savings") or 0
        row = _write_table_row(
            ws4,
            row,
            [
                str(rec.get("priority_num") or ""),
                priority.upper(),
                rec.get("action") or "",
                rec.get("detail") or "",
                f"${rec_savings:,.0f}" if rec_savings > 0 else "-",
            ],
            fonts=[
                f_body,
                pri_font,
                f_body_bold,
                f_body,
                f_green if rec_savings > 0 else f_body,
            ],
            fills=[None, pri_fill, None, None, None],
        )

    row += 2
    ws4.cell(row=row, column=COL_START, value=savings.get("narrative") or "").font = (
        f_footnote
    )

    row += 2
    ws4.cell(
        row=row,
        column=COL_START,
        value=f"Report generated by Nova AI Suite | {datetime.datetime.now().strftime('%Y-%m-%d %H:%M')}",
    ).font = f_footnote

    # Write to bytes
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ═══════════════════════════════════════════════════════════════════════════════
# 9. GENERATE AUDIT PPT
# ═══════════════════════════════════════════════════════════════════════════════


def generate_audit_ppt(
    report_data: Dict[str, Any], client_name: str = "Client"
) -> bytes:
    """Generate 5-slide PPT audit report.

    Slides:
        1. Audit Overview (title + key metrics)
        2. Findings (scorecard + grades)
        3. Savings Opportunity
        4. Missing Channels
        5. Recommendations

    Uses branding: Port Gore #202058, Blue Violet #5A54BD, Downy #6BB3CD.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_SHAPE
    except ImportError:
        logger.error("python-pptx not available for PPT generation")
        return b""

    # Brand colors
    NAVY = RGBColor(0x20, 0x20, 0x58)
    BLUE = RGBColor(0x5A, 0x54, 0xBD)
    TEAL = RGBColor(0x6B, 0xB3, 0xCD)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    OFF_WHITE = RGBColor(0xFF, 0xFD, 0xF9)
    DARK_TEXT = RGBColor(0x20, 0x20, 0x58)
    MUTED_TEXT = RGBColor(0x59, 0x67, 0x80)
    GREEN = RGBColor(0x33, 0x87, 0x21)
    AMBER = RGBColor(0xCE, 0x90, 0x47)
    RED_ACCENT = RGBColor(0xB5, 0x66, 0x9C)
    LIGHT_BG = RGBColor(0xF5, 0xF3, 0xFF)
    LIGHT_GREEN = RGBColor(0xE6, 0xF2, 0xE0)
    LIGHT_RED = RGBColor(0xFD, 0xE8, 0xE8)
    WARM_GRAY = RGBColor(0xEB, 0xE6, 0xE0)

    FONT_TITLE = "Poppins"
    FONT_BODY = "Inter"

    # Widescreen 16:9
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    scorecard = report_data.get("scorecard", {})
    audit_results = report_data.get("audit_results") or []
    missing_channels = report_data.get("missing_channels") or []
    recommendations = report_data.get("recommendations") or []
    savings = report_data.get("savings", {})

    # ── Helpers ─────────────────────────────────────────────────────
    def _add_bg(slide, color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_shape(slide, left, top, width, height, fill_color, line_color=None):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if line_color:
            shape.line.fill.solid()
            shape.line.fill.fore_color.rgb = line_color
            shape.line.width = Pt(1.5)
        else:
            shape.line.fill.background()
        shape.shadow.inherit = False
        return shape

    def _add_text_box(
        slide,
        left,
        top,
        width,
        height,
        text,
        font_name,
        size,
        color,
        bold=False,
        align=PP_ALIGN.LEFT,
    ):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(text)
        p.font.name = font_name
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = align
        return txBox

    # ── SLIDE 1: Audit Overview ──────────────────────────────────────
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_bg(slide1, NAVY)

    _add_text_box(
        slide1,
        Inches(0.8),
        Inches(0.5),
        Inches(11),
        Inches(1),
        "Media Plan Audit Report",
        FONT_TITLE,
        36,
        WHITE,
        bold=True,
    )
    _add_shape(slide1, Inches(0.8), Inches(1.3), Inches(2.5), Inches(0.05), TEAL)
    _add_text_box(
        slide1,
        Inches(0.8),
        Inches(1.6),
        Inches(8),
        Inches(0.5),
        f"{client_name} | {datetime.datetime.now().strftime('%B %d, %Y')}",
        FONT_BODY,
        14,
        TEAL,
    )

    # Key metric cards
    overall_grade = scorecard.get("overall_grade", "N/A")
    efficiency = scorecard.get("budget_efficiency_score") or 0
    coverage = scorecard.get("channel_coverage_score") or 0
    total_savings_val = savings.get("total_savings") or 0

    cards = [
        ("Overall Grade", overall_grade, ""),
        ("Budget Efficiency", f"{efficiency}", "/100"),
        ("Channel Coverage", f"{coverage}", "/100"),
        ("Potential Savings", f"${total_savings_val:,.0f}", ""),
    ]
    for i, (label, value, suffix) in enumerate(cards):
        left = Inches(0.8 + i * 3.1)
        card = _add_shape(
            slide1, left, Inches(2.5), Inches(2.8), Inches(2.0), OFF_WHITE, TEAL
        )

        _add_text_box(
            slide1,
            left + Inches(0.2),
            Inches(2.7),
            Inches(2.4),
            Inches(0.4),
            label,
            FONT_BODY,
            11,
            MUTED_TEXT,
        )
        val_color = (
            GREEN
            if (i == 0 and overall_grade in ("A", "B"))
            or (i == 3 and total_savings_val == 0)
            else (
                AMBER
                if (i == 0 and overall_grade == "C")
                else (
                    RED_ACCENT
                    if (i == 0 and overall_grade in ("D", "F"))
                    else DARK_TEXT
                )
            )
        )
        _add_text_box(
            slide1,
            left + Inches(0.2),
            Inches(3.1),
            Inches(2.4),
            Inches(0.8),
            f"{value}{suffix}",
            FONT_TITLE,
            28,
            val_color,
            bold=True,
            align=PP_ALIGN.CENTER,
        )

    # Industry & channel count
    industry_label = INDUSTRY_LABEL_MAP.get(
        report_data.get("industry") or "",
        (report_data.get("industry") or "").replace("_", " ").title(),
    )
    _add_text_box(
        slide1,
        Inches(0.8),
        Inches(5.0),
        Inches(6),
        Inches(0.4),
        f"Industry: {industry_label}  |  Channels Analyzed: {scorecard.get('channel_count') or 0}  |  "
        f"Missing Channels: {scorecard.get('missing_channel_count') or 0}",
        FONT_BODY,
        11,
        WARM_GRAY,
    )

    _add_shape(
        slide1,
        Inches(0),
        Inches(7.0),
        prs.slide_width,
        Inches(0.5),
        RGBColor(0x15, 0x15, 0x40),
    )
    _add_text_box(
        slide1,
        Inches(0.5),
        Inches(7.05),
        Inches(4),
        Inches(0.4),
        "Powered by Nova AI Suite",
        FONT_BODY,
        10,
        WHITE,
        align=PP_ALIGN.LEFT,
    )

    # ── SLIDE 2: Findings ──────────────────────────────────────────
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide2, OFF_WHITE)

    _add_text_box(
        slide2,
        Inches(0.8),
        Inches(0.4),
        Inches(8),
        Inches(0.7),
        "Audit Findings",
        FONT_TITLE,
        28,
        NAVY,
        bold=True,
    )
    _add_shape(slide2, Inches(0.8), Inches(1.0), Inches(2), Inches(0.04), TEAL)

    # Channel grade cards (up to 8)
    max_channels = min(len(audit_results), 8)
    for i, ar in enumerate(
        sorted(
            audit_results, key=lambda x: x.get("efficiency_score") or 0, reverse=True
        )[:max_channels]
    ):
        col = i % 4
        r = i // 4
        left = Inches(0.5 + col * 3.2)
        top = Inches(1.5 + r * 2.8)

        grade = ar.get("grade", "C")
        card_border = (
            GREEN if grade in ("A", "B") else (AMBER if grade == "C" else RED_ACCENT)
        )
        _add_shape(slide2, left, top, Inches(2.9), Inches(2.4), WHITE, card_border)

        _add_text_box(
            slide2,
            left + Inches(0.15),
            top + Inches(0.1),
            Inches(2.0),
            Inches(0.3),
            ar.get("channel") or "",
            FONT_TITLE,
            12,
            NAVY,
            bold=True,
        )
        _add_text_box(
            slide2,
            left + Inches(2.2),
            top + Inches(0.1),
            Inches(0.5),
            Inches(0.3),
            grade,
            FONT_TITLE,
            18,
            card_border,
            bold=True,
            align=PP_ALIGN.CENTER,
        )

        metrics_lines = [
            f"Budget: ${ar.get('planned_budget') or 0:,.0f}",
            (
                f"CPC: ${ar.get('planned_cpc') or 0:.2f} (Bench: ${ar.get('benchmark_cpc') or 0:.2f})"
                if ar.get("planned_cpc")
                else "CPC: N/A"
            ),
            f"Variance: {ar.get('cpc_variance_pct') or 0:+.0f}%",
            f"Finding: {_finding_text(ar.get('finding') or '')}",
        ]
        for m_idx, m_text in enumerate(metrics_lines):
            _add_text_box(
                slide2,
                left + Inches(0.15),
                top + Inches(0.55 + m_idx * 0.4),
                Inches(2.6),
                Inches(0.35),
                m_text,
                FONT_BODY,
                9,
                MUTED_TEXT if m_idx < 3 else DARK_TEXT,
                bold=(m_idx == 3),
            )

    _add_shape(slide2, Inches(0), Inches(7.0), prs.slide_width, Inches(0.5), NAVY)

    pass  # _finding_text is module-level

    # ── SLIDE 3: Savings Opportunity ──────────────────────────────────
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide3, OFF_WHITE)

    _add_text_box(
        slide3,
        Inches(0.8),
        Inches(0.4),
        Inches(8),
        Inches(0.7),
        "Savings Opportunity",
        FONT_TITLE,
        28,
        NAVY,
        bold=True,
    )
    _add_shape(slide3, Inches(0.8), Inches(1.0), Inches(2), Inches(0.04), TEAL)

    # Big savings number
    if total_savings_val > 0:
        _add_shape(
            slide3,
            Inches(0.8),
            Inches(1.5),
            Inches(11.5),
            Inches(1.2),
            LIGHT_GREEN,
            GREEN,
        )
        _add_text_box(
            slide3,
            Inches(1.2),
            Inches(1.6),
            Inches(10),
            Inches(0.6),
            f"You could save ${total_savings_val:,.0f} by optimizing your media plan",
            FONT_TITLE,
            22,
            GREEN,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        _add_text_box(
            slide3,
            Inches(1.2),
            Inches(2.15),
            Inches(10),
            Inches(0.4),
            f"{savings.get('savings_pct') or 0}% of your ${savings.get('total_budget') or 0:,.0f} total planned budget",
            FONT_BODY,
            12,
            MUTED_TEXT,
            align=PP_ALIGN.CENTER,
        )
    else:
        _add_shape(
            slide3,
            Inches(0.8),
            Inches(1.5),
            Inches(11.5),
            Inches(1.0),
            LIGHT_GREEN,
            GREEN,
        )
        _add_text_box(
            slide3,
            Inches(1.2),
            Inches(1.65),
            Inches(10),
            Inches(0.6),
            "Your planned costs are in line with industry benchmarks",
            FONT_TITLE,
            20,
            GREEN,
            bold=True,
            align=PP_ALIGN.CENTER,
        )

    # Per-channel savings breakdown
    channel_savings = savings.get("channel_savings") or []
    if channel_savings:
        _add_text_box(
            slide3,
            Inches(0.8),
            Inches(3.0),
            Inches(4),
            Inches(0.4),
            "Savings by Channel",
            FONT_TITLE,
            16,
            NAVY,
            bold=True,
        )

        for i, cs in enumerate(channel_savings[:6]):
            top = Inches(3.5 + i * 0.55)
            _add_text_box(
                slide3,
                Inches(1.0),
                top,
                Inches(3.5),
                Inches(0.4),
                cs.get("channel") or "",
                FONT_BODY,
                11,
                DARK_TEXT,
                bold=True,
            )
            _add_text_box(
                slide3,
                Inches(5.0),
                top,
                Inches(2.5),
                Inches(0.4),
                f"${cs.get('savings') or 0:,.0f}",
                FONT_BODY,
                11,
                GREEN,
                bold=True,
            )
            _add_text_box(
                slide3,
                Inches(7.5),
                top,
                Inches(2),
                Inches(0.4),
                f"({cs.get('savings_pct') or 0:.0f}% of channel budget)",
                FONT_BODY,
                9,
                MUTED_TEXT,
            )

    _add_shape(slide3, Inches(0), Inches(7.0), prs.slide_width, Inches(0.5), NAVY)

    # ── SLIDE 4: Missing Channels ──────────────────────────────────
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide4, OFF_WHITE)

    _add_text_box(
        slide4,
        Inches(0.8),
        Inches(0.4),
        Inches(8),
        Inches(0.7),
        "Missing Channels",
        FONT_TITLE,
        28,
        NAVY,
        bold=True,
    )
    _add_shape(slide4, Inches(0.8), Inches(1.0), Inches(2), Inches(0.04), TEAL)

    if missing_channels:
        _add_text_box(
            slide4,
            Inches(0.8),
            Inches(1.3),
            Inches(10),
            Inches(0.4),
            f"{len(missing_channels)} recommended channel(s) not in your current plan",
            FONT_BODY,
            12,
            MUTED_TEXT,
        )

        max_cards = min(len(missing_channels), 6)
        for i, mc in enumerate(missing_channels[:max_cards]):
            col = i % 3
            r = i // 3
            left = Inches(0.5 + col * 4.2)
            top = Inches(1.9 + r * 2.5)

            importance = mc.get("importance", "medium")
            if importance == "critical":
                border_color = RED_ACCENT
                badge_text = "CRITICAL"
            elif importance == "high":
                border_color = AMBER
                badge_text = "HIGH"
            else:
                border_color = TEAL
                badge_text = "MEDIUM"

            _add_shape(slide4, left, top, Inches(3.8), Inches(2.1), WHITE, border_color)

            # Badge
            badge = _add_shape(
                slide4,
                left + Inches(0.15),
                top + Inches(0.15),
                Inches(1.2),
                Inches(0.3),
                border_color,
            )
            _add_text_box(
                slide4,
                left + Inches(0.15),
                top + Inches(0.16),
                Inches(1.2),
                Inches(0.3),
                badge_text,
                FONT_BODY,
                8,
                WHITE,
                bold=True,
                align=PP_ALIGN.CENTER,
            )

            # Channel name
            _add_text_box(
                slide4,
                left + Inches(0.15),
                top + Inches(0.55),
                Inches(3.4),
                Inches(0.3),
                mc.get("channel") or "",
                FONT_TITLE,
                13,
                NAVY,
                bold=True,
            )

            # Reason
            _add_text_box(
                slide4,
                left + Inches(0.15),
                top + Inches(0.9),
                Inches(3.4),
                Inches(1.0),
                mc.get("reason") or "",
                FONT_BODY,
                9,
                MUTED_TEXT,
            )
    else:
        _add_shape(
            slide4, Inches(2), Inches(2.5), Inches(9), Inches(1.5), LIGHT_GREEN, GREEN
        )
        _add_text_box(
            slide4,
            Inches(2.5),
            Inches(2.9),
            Inches(8),
            Inches(0.8),
            "Your media plan covers all recommended channels.\nExcellent channel coverage!",
            FONT_TITLE,
            18,
            GREEN,
            bold=True,
            align=PP_ALIGN.CENTER,
        )

    _add_shape(slide4, Inches(0), Inches(7.0), prs.slide_width, Inches(0.5), NAVY)

    # ── SLIDE 5: Recommendations ──────────────────────────────────
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide5, OFF_WHITE)

    _add_text_box(
        slide5,
        Inches(0.8),
        Inches(0.4),
        Inches(8),
        Inches(0.7),
        "Top Recommendations",
        FONT_TITLE,
        28,
        NAVY,
        bold=True,
    )
    _add_shape(slide5, Inches(0.8), Inches(1.0), Inches(2), Inches(0.04), TEAL)

    max_recs = min(len(recommendations), 6)
    for i, rec in enumerate(recommendations[:max_recs]):
        priority = rec.get("priority", "medium")
        top = Inches(1.4 + i * 0.9)

        if priority == "high":
            action_color = RED_ACCENT
            bg_color = LIGHT_RED
        elif priority == "medium":
            action_color = AMBER
            bg_color = RGBColor(0xFE, 0xF3, 0xC7)
        else:
            action_color = TEAL
            bg_color = RGBColor(0xE0, 0xF2, 0xFE)

        _add_shape(
            slide5, Inches(0.5), top, Inches(12.3), Inches(0.75), bg_color, action_color
        )

        # Priority badge
        badge = _add_shape(
            slide5,
            Inches(0.7),
            top + Inches(0.12),
            Inches(1.2),
            Inches(0.45),
            action_color,
        )
        _add_text_box(
            slide5,
            Inches(0.7),
            top + Inches(0.15),
            Inches(1.2),
            Inches(0.4),
            priority.upper(),
            FONT_BODY,
            9,
            WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )

        # Number
        _add_text_box(
            slide5,
            Inches(2.1),
            top + Inches(0.08),
            Inches(0.5),
            Inches(0.35),
            f"#{rec.get('priority_num', i + 1)}",
            FONT_TITLE,
            12,
            NAVY,
            bold=True,
        )

        # Action
        _add_text_box(
            slide5,
            Inches(2.7),
            top + Inches(0.08),
            Inches(5),
            Inches(0.35),
            rec.get("action") or "",
            FONT_TITLE,
            11,
            NAVY,
            bold=True,
        )

        # Savings if any
        rec_savings = rec.get("savings") or 0
        if rec_savings > 0:
            _add_text_box(
                slide5,
                Inches(10.5),
                top + Inches(0.15),
                Inches(2),
                Inches(0.35),
                f"Save ${rec_savings:,.0f}",
                FONT_BODY,
                10,
                GREEN,
                bold=True,
                align=PP_ALIGN.RIGHT,
            )

    # Narrative
    _add_text_box(
        slide5,
        Inches(0.8),
        Inches(6.3),
        Inches(11.5),
        Inches(0.5),
        savings.get("narrative") or "",
        FONT_BODY,
        10,
        TEAL,
        bold=True,
    )

    _add_shape(slide5, Inches(0), Inches(7.0), prs.slide_width, Inches(0.5), NAVY)
    _add_text_box(
        slide5,
        Inches(0.5),
        Inches(7.05),
        Inches(4),
        Inches(0.4),
        "Powered by Nova AI Suite",
        FONT_BODY,
        10,
        WHITE,
        align=PP_ALIGN.LEFT,
    )

    # Write to bytes
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ═══════════════════════════════════════════════════════════════════════════════
# 10. ORCHESTRATOR -- single entry point
# ═══════════════════════════════════════════════════════════════════════════════


def run_full_audit(
    file_bytes: bytes,
    filename: str,
    industry: str = "general_entry_level",
    roles: Optional[List[str]] = None,
    locations: Optional[List[str]] = None,
    client_name: str = "Client",
) -> Dict[str, Any]:
    """Full audit pipeline: parse -> benchmark -> audit -> missing -> savings -> recommend -> scorecard.

    Thread-safe, never raises. All errors return structured error dicts.

    Returns a complete report_data dict ready for Excel/PPT generation, including:
        success, audit_results, missing_channels, savings, recommendations,
        scorecard, line_items, industry, generated_at
    """
    try:
        # 1. Parse the media plan
        line_items = parse_media_plan(file_bytes, filename)
        if not line_items:
            return {
                "error": (
                    "Could not parse the media plan. Please ensure your file has columns for "
                    "Channel/Platform and Budget/Spend at minimum. Supported formats: Excel (.xlsx/.xls) and CSV."
                ),
                "success": False,
            }

        # 2. Determine collar type
        collar_type = "mixed"
        if roles and _HAS_COLLAR_INTEL:
            try:
                result = _collar_intel.classify_collar(roles[0], industry=industry)
                collar_type = result.get("collar_type", "mixed")
            except Exception:
                pass

        # 3. Get all benchmarks (cache for reuse across line items)
        location = ""
        if locations:
            location = (
                locations[0] if isinstance(locations[0], str) else str(locations[0])
            )
        benchmarks_cache = _get_all_benchmarks(industry, collar_type, location)

        # 4. Audit each line item
        audit_results: List[Dict[str, Any]] = []
        for item in line_items:
            audit = audit_line_item(
                item,
                industry=industry,
                roles=roles,
                locations=locations,
                collar_type=collar_type,
                benchmarks_cache=benchmarks_cache,
            )
            audit_results.append(audit)

        # 5. Identify missing channels
        current_channels = [item.get("channel") or "" for item in line_items]
        missing_channels = identify_missing_channels(
            current_channels, industry, roles, collar_type
        )

        # 6. Calculate savings
        savings_data = calculate_savings_potential(audit_results)

        # 7. Generate recommendations
        recs = generate_recommendations(audit_results, missing_channels, industry)

        # 8. Generate scorecard
        scorecard = generate_audit_scorecard(audit_results, missing_channels)

        return {
            "success": True,
            "client_name": client_name or "Client",
            "industry": industry,
            "industry_label": INDUSTRY_LABEL_MAP.get(
                industry, industry.replace("_", " ").title()
            ),
            "collar_type": collar_type,
            "line_items": line_items,
            "audit_results": audit_results,
            "missing_channels": missing_channels,
            "savings": savings_data,
            "recommendations": recs,
            "scorecard": scorecard,
            "channel_count": len(line_items),
            "total_planned_budget": sum(item.get("budget") or 0 for item in line_items),
            "generated_at": datetime.datetime.now().isoformat(),
        }

    except Exception as exc:
        logger.exception("Media plan audit failed: %s", exc)
        return {"error": f"Audit failed: {str(exc)}", "success": False}
