"""
Budget Simulator -- Multi-Scenario Budget Allocation Simulator

Interactive backend for the budget simulator tool. Users adjust budget sliders
across 10 recruitment channels and see projected outcomes in real-time.
Supports saving and comparing up to 3 scenarios side-by-side.

Channels:
    programmatic_dsp, job_boards, paid_search, paid_social, linkedin,
    career_site_seo, employee_referrals, staffing_agency,
    events_career_fairs, niche_diversity_boards

All simulation uses cached benchmarks only (no external API calls).
Thread-safe, fast (<50ms per simulation).
"""

from __future__ import annotations

import io
import logging
import math
import datetime
import threading
from typing import Any, Dict, List, Optional, Tuple

logger = logging.getLogger(__name__)

# ── Lazy imports (try/except pattern matching codebase conventions) ──

try:
    import budget_engine as _budget_engine
    _HAS_BUDGET_ENGINE = True
except ImportError:
    _budget_engine = None  # type: ignore
    _HAS_BUDGET_ENGINE = False

try:
    import trend_engine as _trend_engine
    _HAS_TREND_ENGINE = True
except ImportError:
    _trend_engine = None  # type: ignore
    _HAS_TREND_ENGINE = False

try:
    import collar_intelligence as _collar_intel
    _HAS_COLLAR_INTEL = True
except ImportError:
    _collar_intel = None  # type: ignore
    _HAS_COLLAR_INTEL = False

try:
    from shared_utils import INDUSTRY_LABEL_MAP, parse_budget
except ImportError:
    INDUSTRY_LABEL_MAP = {}
    def parse_budget(v, *, default=100_000.0):  # type: ignore
        try:
            return float(v)
        except (TypeError, ValueError):
            return default

# ═══════════════════════════════════════════════════════════════════════════════
# CONSTANTS
# ═══════════════════════════════════════════════════════════════════════════════

# Thread lock for any shared mutable state (currently none, but future-proof)
_LOCK = threading.Lock()

# Canonical simulator channels with display names, icons, and platform mappings
SIMULATOR_CHANNELS: Dict[str, Dict[str, Any]] = {
    "programmatic_dsp": {
        "label": "Programmatic & DSP",
        "icon": "📡",
        "category": "programmatic",
        "platform": "programmatic",
        "description": "Automated programmatic job advertising across ad exchanges",
    },
    "job_boards": {
        "label": "Job Boards",
        "icon": "📋",
        "category": "job_board",
        "platform": "indeed",
        "description": "Major job boards (Indeed, ZipRecruiter, Monster)",
    },
    "paid_search": {
        "label": "Paid Search / SEM",
        "icon": "🔍",
        "category": "search",
        "platform": "google_search",
        "description": "Google Ads, Bing Ads for recruitment keywords",
    },
    "paid_social": {
        "label": "Paid Social",
        "icon": "📱",
        "category": "social",
        "platform": "meta_facebook",
        "description": "Meta (Facebook/Instagram), TikTok recruitment ads",
    },
    "linkedin": {
        "label": "LinkedIn",
        "icon": "💼",
        "category": "niche_board",
        "platform": "linkedin",
        "description": "LinkedIn Recruiter, Sponsored Jobs, InMail campaigns",
    },
    "career_site_seo": {
        "label": "Career Site & SEO",
        "icon": "🌐",
        "category": "career_site",
        "platform": "google_search",
        "description": "Organic career site optimization, SEO, content marketing",
    },
    "employee_referrals": {
        "label": "Employee Referrals",
        "icon": "🤝",
        "category": "referral",
        "platform": None,
        "description": "Internal referral programs with bonuses",
    },
    "staffing_agency": {
        "label": "Staffing Agency",
        "icon": "🏢",
        "category": "staffing",
        "platform": None,
        "description": "External staffing and recruitment agencies",
    },
    "events_career_fairs": {
        "label": "Events & Career Fairs",
        "icon": "🎪",
        "category": "events",
        "platform": None,
        "description": "Job fairs, campus recruiting, industry events",
    },
    "niche_diversity_boards": {
        "label": "Niche & Diversity Boards",
        "icon": "🌈",
        "category": "niche_board",
        "platform": "linkedin",
        "description": "Specialized boards (Dice, BuiltIn, DiversityJobs, etc.)",
    },
}

CHANNEL_KEYS = list(SIMULATOR_CHANNELS.keys())

# Industry key mapping from user-facing labels to internal keys
SIMULATOR_INDUSTRY_MAP: Dict[str, str] = {
    "Technology": "tech_engineering",
    "Healthcare": "healthcare_medical",
    "Finance": "finance_banking",
    "Manufacturing": "automotive",
    "Retail": "retail_consumer",
    "Transportation": "logistics_supply_chain",
    "Hospitality": "hospitality_travel",
    "Construction": "construction_real_estate",
    "Government": "general_entry_level",
    "Education": "education",
    "Energy": "energy_utilities",
    "Agriculture": "blue_collar_trades",
    "Telecommunications": "telecommunications",
    "Media": "media_entertainment",
    "Real Estate": "construction_real_estate",
    "Logistics": "logistics_supply_chain",
    "Nonprofit": "general_entry_level",
}

SIMULATOR_INDUSTRIES = list(SIMULATOR_INDUSTRY_MAP.keys())

# Default channel allocations by collar type
_DEFAULT_ALLOC_BLUE: Dict[str, float] = {
    "programmatic_dsp": 30,
    "job_boards": 25,
    "paid_search": 5,
    "paid_social": 15,
    "linkedin": 3,
    "career_site_seo": 5,
    "employee_referrals": 7,
    "staffing_agency": 2,
    "events_career_fairs": 5,
    "niche_diversity_boards": 3,
}

_DEFAULT_ALLOC_WHITE: Dict[str, float] = {
    "programmatic_dsp": 15,
    "job_boards": 12,
    "paid_search": 8,
    "paid_social": 8,
    "linkedin": 25,
    "career_site_seo": 8,
    "employee_referrals": 10,
    "staffing_agency": 3,
    "events_career_fairs": 3,
    "niche_diversity_boards": 8,
}

_DEFAULT_ALLOC_GREY: Dict[str, float] = {
    "programmatic_dsp": 20,
    "job_boards": 18,
    "paid_search": 5,
    "paid_social": 10,
    "linkedin": 10,
    "career_site_seo": 7,
    "employee_referrals": 8,
    "staffing_agency": 5,
    "events_career_fairs": 5,
    "niche_diversity_boards": 12,
}

_DEFAULT_ALLOC_MIXED: Dict[str, float] = {
    "programmatic_dsp": 22,
    "job_boards": 18,
    "paid_search": 7,
    "paid_social": 12,
    "linkedin": 12,
    "career_site_seo": 7,
    "employee_referrals": 8,
    "staffing_agency": 3,
    "events_career_fairs": 4,
    "niche_diversity_boards": 7,
}

# Fallback CPC benchmarks per channel (USD) -- used when trend_engine unavailable
_FALLBACK_CPC: Dict[str, float] = {
    "programmatic_dsp": 0.65,
    "job_boards": 0.85,
    "paid_search": 2.50,
    "paid_social": 1.20,
    "linkedin": 3.80,
    "career_site_seo": 0.30,
    "employee_referrals": 0.00,
    "staffing_agency": 0.00,
    "events_career_fairs": 0.00,
    "niche_diversity_boards": 1.40,
}

# Fallback apply rates per channel
_FALLBACK_APPLY_RATE: Dict[str, float] = {
    "programmatic_dsp": 0.06,
    "job_boards": 0.08,
    "paid_search": 0.05,
    "paid_social": 0.03,
    "linkedin": 0.04,
    "career_site_seo": 0.12,
    "employee_referrals": 0.25,
    "staffing_agency": 0.20,
    "events_career_fairs": 0.15,
    "niche_diversity_boards": 0.10,
}

# Fallback hire rate (applications -> hires)
_FALLBACK_HIRE_RATE: float = 0.02


# ═══════════════════════════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════════════════════════

def _safe_div(num: float, den: float, default: float = 0.0) -> float:
    """Division that never raises ZeroDivisionError."""
    return num / den if den != 0 else default


def _resolve_industry_key(industry: str) -> str:
    """Map a user-facing industry label to an internal industry key."""
    if not industry:
        return "general_entry_level"
    # Direct match in SIMULATOR_INDUSTRY_MAP
    if industry in SIMULATOR_INDUSTRY_MAP:
        return SIMULATOR_INDUSTRY_MAP[industry]
    # Already an internal key
    if industry in (INDUSTRY_LABEL_MAP or {}):
        return industry
    # Fuzzy match
    industry_lower = industry.lower().strip()
    for label, key in SIMULATOR_INDUSTRY_MAP.items():
        if label.lower() == industry_lower:
            return key
    for key in (INDUSTRY_LABEL_MAP or {}):
        if industry_lower in key or key in industry_lower:
            return key
    return "general_entry_level"


def _resolve_collar_type(roles: str, industry: str) -> str:
    """Determine collar type from role text and industry."""
    if not roles:
        # Fall back to industry-based collar
        if _HAS_COLLAR_INTEL:
            ind_key = _resolve_industry_key(industry)
            default_collar = getattr(_collar_intel, '_INDUSTRY_DEFAULT_COLLAR', {})
            return default_collar.get(ind_key, "white_collar")
        return "white_collar"

    if _HAS_COLLAR_INTEL:
        try:
            result = _collar_intel.classify_collar(
                role=roles,
                industry=_resolve_industry_key(industry),
            )
            return result.get("collar_type", "white_collar")
        except Exception:
            pass

    # Keyword fallback
    roles_lower = roles.lower()
    blue_kw = {"driver", "warehouse", "forklift", "construction", "mechanic",
               "welder", "laborer", "factory", "production", "cook", "cleaner"}
    grey_kw = {"nurse", "therapist", "technician", "emt", "paramedic", "medical"}
    if any(kw in roles_lower for kw in blue_kw):
        return "blue_collar"
    if any(kw in roles_lower for kw in grey_kw):
        return "grey_collar"
    return "white_collar"


def _get_channel_cpc(channel_key: str, industry_key: str, collar_type: str,
                     location: str = "") -> float:
    """Get CPC for a channel using trend_engine if available, else fallback."""
    ch_info = SIMULATOR_CHANNELS.get(channel_key, {})
    platform = ch_info.get("platform")

    if _HAS_TREND_ENGINE and platform:
        try:
            month = datetime.datetime.now().month
            result = _trend_engine.get_benchmark(
                platform=platform,
                industry=industry_key,
                metric="cpc",
                collar_type=collar_type,
                location=location,
                month=month,
            )
            if result and isinstance(result, dict):
                val = result.get("value", 0)
                if isinstance(val, (int, float)) and val > 0:
                    return float(val)
        except Exception:
            pass

    return _FALLBACK_CPC.get(channel_key, 0.85)


def _get_channel_apply_rate(channel_key: str, collar_type: str) -> float:
    """Get apply rate for a channel, adjusted by collar type."""
    base_rate = _FALLBACK_APPLY_RATE.get(channel_key, 0.05)

    if _HAS_BUDGET_ENGINE:
        category = SIMULATOR_CHANNELS.get(channel_key, {}).get("category", "job_board")
        try:
            collar_mult = _budget_engine._get_collar_apply_rate_adjustment(category, collar_type)
            return round(base_rate * collar_mult, 4)
        except Exception:
            pass

    # Simple collar adjustments without budget_engine
    if collar_type == "blue_collar":
        if channel_key in ("job_boards", "programmatic_dsp", "paid_social"):
            return base_rate * 1.3
        if channel_key == "linkedin":
            return base_rate * 0.7
    elif collar_type == "white_collar":
        if channel_key == "linkedin":
            return base_rate * 1.4
        if channel_key in ("job_boards", "programmatic_dsp"):
            return base_rate * 0.9

    return base_rate


def _get_hire_rate(collar_type: str) -> float:
    """Get hire rate (applications -> hires) based on collar type."""
    if _HAS_BUDGET_ENGINE:
        tier_map = {
            "blue_collar": "Hourly / Entry-Level",
            "white_collar": "Professional / White-Collar",
            "grey_collar": "Clinical / Licensed",
            "pink_collar": "Hourly / Entry-Level",
        }
        tier = tier_map.get(collar_type, "default")
        return _budget_engine.HIRE_RATE_BY_TIER.get(tier, _FALLBACK_HIRE_RATE)

    # Fallback rates
    return {
        "blue_collar": 0.06,
        "white_collar": 0.02,
        "grey_collar": 0.03,
        "pink_collar": 0.05,
    }.get(collar_type, _FALLBACK_HIRE_RATE)


def _score_roi(cost_per_hire: float, industry_key: str) -> int:
    """Score ROI on a 1-10 scale based on industry average CPH."""
    if _HAS_BUDGET_ENGINE:
        try:
            avg_cph = _budget_engine._industry_avg_cph(industry_key)
        except Exception:
            avg_cph = 6000.0
    else:
        avg_cph = 6000.0

    if avg_cph <= 0 or cost_per_hire <= 0:
        return 5
    ratio = cost_per_hire / avg_cph
    score = 10 - (ratio - 0.2) * (9 / 2.8)
    return int(max(1, min(10, round(score))))


# ═══════════════════════════════════════════════════════════════════════════════
# CORE PUBLIC FUNCTIONS
# ═══════════════════════════════════════════════════════════════════════════════

def get_default_allocation(
    total_budget: float,
    industry: str = "",
    roles: str = "",
    locations: str = "",
) -> Dict[str, Any]:
    """Return default % allocation per channel with projected metrics.

    Uses collar_intelligence and budget_engine when available to produce
    industry/role-appropriate defaults.

    Args:
        total_budget: Total campaign budget in USD.
        industry: Industry label (e.g. "Technology", "Healthcare").
        roles: Free-text role description (e.g. "Software Engineer").
        locations: Free-text location (e.g. "San Francisco, CA").

    Returns:
        Dict with:
            - channel_allocations: {channel_key: percentage}
            - projected_metrics: per-channel metric projections
            - metadata: industry_key, collar_type, total_budget
    """
    industry_key = _resolve_industry_key(industry)
    collar_type = _resolve_collar_type(roles, industry)

    # Select default allocation based on collar type
    collar_alloc_map = {
        "blue_collar": _DEFAULT_ALLOC_BLUE,
        "white_collar": _DEFAULT_ALLOC_WHITE,
        "grey_collar": _DEFAULT_ALLOC_GREY,
        "pink_collar": _DEFAULT_ALLOC_BLUE,  # similar distribution to blue
    }
    base_alloc = collar_alloc_map.get(collar_type, _DEFAULT_ALLOC_MIXED)

    # If collar_intelligence has a blended allocation, use it for guidance
    if _HAS_COLLAR_INTEL and roles:
        try:
            blend_result = _collar_intel.get_blended_allocation(
                [{"role": roles, "count": 1, "industry": industry_key}]
            )
            blended_mix = blend_result.get("blended_channel_mix", {})
            if blended_mix:
                # Map collar_intelligence channel keys to simulator channel keys
                _ci_to_sim = {
                    "programmatic": "programmatic_dsp",
                    "global_job_boards": "job_boards",
                    "social_media": "paid_social",
                    "linkedin": "linkedin",
                    "niche_boards": "niche_diversity_boards",
                    "employer_branding": "career_site_seo",
                    "regional_local": "events_career_fairs",
                    "search": "paid_search",
                }
                mapped: Dict[str, float] = {}
                for ci_key, pct in blended_mix.items():
                    sim_key = _ci_to_sim.get(ci_key)
                    if sim_key:
                        mapped[sim_key] = round(pct * 100, 1)

                # Fill in unmapped channels from base
                for ch_key in CHANNEL_KEYS:
                    if ch_key not in mapped:
                        mapped[ch_key] = base_alloc.get(ch_key, 5.0)

                # Normalize to 100%
                total_pct = sum(mapped.values())
                if total_pct > 0:
                    base_alloc = {k: round(v / total_pct * 100, 1) for k, v in mapped.items()}
        except Exception as e:
            logger.debug("collar_intelligence blended allocation failed: %s", e)

    # Ensure all channels present and sum to 100
    channel_allocations: Dict[str, float] = {}
    for ch_key in CHANNEL_KEYS:
        channel_allocations[ch_key] = base_alloc.get(ch_key, 5.0)

    total_pct = sum(channel_allocations.values())
    if abs(total_pct - 100.0) > 0.1:
        factor = 100.0 / total_pct
        channel_allocations = {k: round(v * factor, 1) for k, v in channel_allocations.items()}
        # Fix rounding residual
        diff = 100.0 - sum(channel_allocations.values())
        if abs(diff) > 0.01:
            max_ch = max(channel_allocations, key=channel_allocations.get)
            channel_allocations[max_ch] = round(channel_allocations[max_ch] + diff, 1)

    return {
        "channel_allocations": channel_allocations,
        "total_budget": total_budget,
        "metadata": {
            "industry": industry,
            "industry_key": industry_key,
            "collar_type": collar_type,
            "roles": roles,
            "locations": locations,
        },
    }


def simulate_scenario(
    channel_allocations: Dict[str, float],
    total_budget: float,
    industry: str = "",
    roles: str = "",
    locations: str = "",
) -> Dict[str, Any]:
    """Simulate a scenario given channel % allocations and total budget.

    Computes for each channel: dollar_amount, projected_clicks,
    projected_applies, projected_hires, cpc, cpa, cph, roi_score.

    Args:
        channel_allocations: {channel_key: percentage} (should sum to 100).
        total_budget: Total campaign budget in USD.
        industry: Industry label or key.
        roles: Free-text role description.
        locations: Free-text location.

    Returns:
        Full scenario dict with per-channel and aggregate metrics.
    """
    total_budget = max(0.0, total_budget)
    industry_key = _resolve_industry_key(industry)
    collar_type = _resolve_collar_type(roles, industry)
    hire_rate = _get_hire_rate(collar_type)

    # Normalize allocations to sum to 100
    pct_sum = sum(channel_allocations.values())
    if pct_sum <= 0:
        pct_sum = 100.0
    norm = 100.0 / pct_sum

    channels: Dict[str, Dict[str, Any]] = {}
    total_clicks = 0
    total_applies = 0
    total_hires = 0
    total_spend = 0.0

    for ch_key in CHANNEL_KEYS:
        raw_pct = channel_allocations.get(ch_key, 0.0)
        pct = raw_pct * norm
        dollars = round(total_budget * pct / 100.0, 2)
        total_spend += dollars

        cpc = _get_channel_cpc(ch_key, industry_key, collar_type, locations)
        apply_rate = _get_channel_apply_rate(ch_key, collar_type)

        if cpc > 0:
            clicks = max(0, int(dollars / cpc))
            applies = max(0, int(clicks * apply_rate))
            hires = max(0, int(applies * hire_rate))
        else:
            # Flat-cost channels (referrals, events, staffing)
            clicks = 0
            applies = max(1, int(dollars / 50.0)) if dollars > 0 else 0
            hires = max(0, int(applies * hire_rate * 2))  # higher quality

        cpa = _safe_div(dollars, max(applies, 1), dollars)
        cph = _safe_div(dollars, max(hires, 1), dollars)
        roi = _score_roi(cph, industry_key)

        total_clicks += clicks
        total_applies += applies
        total_hires += hires

        channels[ch_key] = {
            "channel_key": ch_key,
            "label": SIMULATOR_CHANNELS[ch_key]["label"],
            "icon": SIMULATOR_CHANNELS[ch_key]["icon"],
            "percentage": round(pct, 1),
            "dollar_amount": dollars,
            "projected_clicks": clicks,
            "projected_applies": applies,
            "projected_hires": hires,
            "cpc": round(cpc, 2),
            "cpa": round(cpa, 2),
            "cph": round(cph, 2),
            "roi_score": roi,
        }

    blended_cpa = _safe_div(total_spend, max(total_applies, 1), 0.0)
    blended_cpc = _safe_div(total_spend, max(total_clicks, 1), 0.0)
    blended_cph = _safe_div(total_spend, max(total_hires, 1), 0.0)
    overall_roi = _score_roi(blended_cph, industry_key)

    return {
        "channels": channels,
        "summary": {
            "total_budget": total_budget,
            "total_clicks": total_clicks,
            "total_applies": total_applies,
            "total_hires": total_hires,
            "blended_cpc": round(blended_cpc, 2),
            "blended_cpa": round(blended_cpa, 2),
            "blended_cph": round(blended_cph, 2),
            "roi_score": overall_roi,
        },
        "metadata": {
            "industry": industry,
            "industry_key": industry_key,
            "collar_type": collar_type,
            "hire_rate": hire_rate,
            "roles": roles,
            "locations": locations,
            "timestamp": datetime.datetime.utcnow().isoformat() + "Z",
        },
    }


def compare_scenarios(scenarios_list: List[Dict[str, Any]]) -> Dict[str, Any]:
    """Compare up to 3 scenario dicts side-by-side.

    Args:
        scenarios_list: List of scenario dicts (output of simulate_scenario).

    Returns:
        Comparison dict with: total_hires, total_applies, blended_cpa,
        blended_cpc, best_scenario_index, per-channel deltas.
    """
    if not scenarios_list:
        return {"error": "No scenarios provided", "scenarios": [], "best_scenario_index": -1}

    scenario_summaries = []
    for i, sc in enumerate(scenarios_list):
        summary = sc.get("summary", {})
        scenario_summaries.append({
            "index": i,
            "label": f"Scenario {i + 1}",
            "total_budget": summary.get("total_budget", 0),
            "total_clicks": summary.get("total_clicks", 0),
            "total_applies": summary.get("total_applies", 0),
            "total_hires": summary.get("total_hires", 0),
            "blended_cpc": summary.get("blended_cpc", 0),
            "blended_cpa": summary.get("blended_cpa", 0),
            "blended_cph": summary.get("blended_cph", 0),
            "roi_score": summary.get("roi_score", 0),
        })

    # Determine best scenario based on composite score:
    # maximize hires, minimize CPA
    best_idx = 0
    best_score = -float("inf")
    for i, s in enumerate(scenario_summaries):
        # Composite: weight hires positively, CPA negatively
        hires = s["total_hires"]
        cpa = s["blended_cpa"]
        score = hires * 100 - cpa  # simple composite
        if score > best_score:
            best_score = score
            best_idx = i

    # Per-channel deltas (compare each scenario's channel % to scenario 1)
    channel_deltas: Dict[str, List[Dict[str, float]]] = {}
    base_channels = scenarios_list[0].get("channels", {}) if scenarios_list else {}

    for ch_key in CHANNEL_KEYS:
        deltas = []
        base_pct = base_channels.get(ch_key, {}).get("percentage", 0)
        base_dollars = base_channels.get(ch_key, {}).get("dollar_amount", 0)

        for i, sc in enumerate(scenarios_list):
            sc_ch = sc.get("channels", {}).get(ch_key, {})
            sc_pct = sc_ch.get("percentage", 0)
            sc_dollars = sc_ch.get("dollar_amount", 0)
            deltas.append({
                "scenario_index": i,
                "percentage": sc_pct,
                "dollar_amount": sc_dollars,
                "pct_delta": round(sc_pct - base_pct, 1) if i > 0 else 0,
                "dollar_delta": round(sc_dollars - base_dollars, 2) if i > 0 else 0,
                "applies": sc_ch.get("projected_applies", 0),
                "hires": sc_ch.get("projected_hires", 0),
            })
        channel_deltas[ch_key] = deltas

    # Identify best/worst per metric
    metrics_best_worst: Dict[str, Dict[str, int]] = {}
    for metric in ("total_clicks", "total_applies", "total_hires", "blended_cpa", "blended_cpc", "roi_score"):
        values = [s.get(metric, 0) for s in scenario_summaries]
        # For CPA/CPC lower is better; for others higher is better
        if metric in ("blended_cpa", "blended_cpc", "blended_cph"):
            non_zero = [v for v in values if v > 0]
            best = values.index(min(non_zero)) if non_zero else 0
            worst = values.index(max(values)) if values else 0
        else:
            best = values.index(max(values)) if values else 0
            worst = values.index(min(values)) if values else 0
        metrics_best_worst[metric] = {"best": best, "worst": worst}

    return {
        "scenarios": scenario_summaries,
        "best_scenario_index": best_idx,
        "channel_deltas": channel_deltas,
        "metrics_best_worst": metrics_best_worst,
        "recommendation": (
            f"Scenario {best_idx + 1} delivers the best overall outcome "
            f"with {scenario_summaries[best_idx]['total_hires']} projected hires "
            f"at ${scenario_summaries[best_idx]['blended_cpa']:.2f} CPA."
        ) if scenario_summaries else "",
    }


def optimize_for_goal(
    total_budget: float,
    goal: str = "balanced",
    industry: str = "",
    roles: str = "",
    locations: str = "",
) -> Dict[str, Any]:
    """Return optimized channel allocation for a specified goal.

    Goals:
        - minimize_cpa: Favor high apply-rate channels
        - maximize_hires: Favor channels with best hire efficiency
        - maximize_applies: Favor channels with highest volume
        - balanced: Weighted mix across all metrics

    Args:
        total_budget: Total budget in USD.
        goal: Optimization goal string.
        industry: Industry label.
        roles: Role description.
        locations: Location.

    Returns:
        Dict with optimized channel_allocations and projected scenario.
    """
    industry_key = _resolve_industry_key(industry)
    collar_type = _resolve_collar_type(roles, industry)
    hire_rate = _get_hire_rate(collar_type)

    # Step 1: Compute efficiency scores per channel
    channel_scores: Dict[str, float] = {}

    for ch_key in CHANNEL_KEYS:
        cpc = _get_channel_cpc(ch_key, industry_key, collar_type, locations)
        apply_rate = _get_channel_apply_rate(ch_key, collar_type)

        if cpc > 0:
            clicks_per_dollar = 1.0 / cpc
            applies_per_dollar = clicks_per_dollar * apply_rate
            hires_per_dollar = applies_per_dollar * hire_rate
            cpa_efficiency = applies_per_dollar  # higher = lower CPA
        else:
            # Flat-cost channels
            applies_per_dollar = 1.0 / 50.0  # heuristic
            hires_per_dollar = applies_per_dollar * hire_rate * 2
            cpa_efficiency = applies_per_dollar
            clicks_per_dollar = 0.0

        if goal == "minimize_cpa":
            channel_scores[ch_key] = cpa_efficiency
        elif goal == "maximize_hires":
            channel_scores[ch_key] = hires_per_dollar
        elif goal == "maximize_applies":
            channel_scores[ch_key] = applies_per_dollar
        else:  # balanced
            channel_scores[ch_key] = (
                cpa_efficiency * 0.3 +
                hires_per_dollar * 1000 * 0.4 +
                applies_per_dollar * 0.3
            )

    # Step 2: Convert scores to allocations
    # Ensure minimum 2% for each channel (floor), distribute rest by score
    min_pct = 2.0
    floor_total = min_pct * len(CHANNEL_KEYS)  # 20%
    remaining = 100.0 - floor_total  # 80% to distribute by score

    total_score = sum(channel_scores.values())
    if total_score <= 0:
        total_score = 1.0

    optimized_alloc: Dict[str, float] = {}
    for ch_key in CHANNEL_KEYS:
        score_share = (channel_scores[ch_key] / total_score) * remaining
        optimized_alloc[ch_key] = round(min_pct + score_share, 1)

    # Normalize to exactly 100%
    alloc_sum = sum(optimized_alloc.values())
    if abs(alloc_sum - 100.0) > 0.05:
        factor = 100.0 / alloc_sum
        optimized_alloc = {k: round(v * factor, 1) for k, v in optimized_alloc.items()}
        diff = 100.0 - sum(optimized_alloc.values())
        if abs(diff) > 0.01:
            max_ch = max(optimized_alloc, key=optimized_alloc.get)
            optimized_alloc[max_ch] = round(optimized_alloc[max_ch] + diff, 1)

    # Step 3: Run simulation with optimized allocation
    scenario = simulate_scenario(optimized_alloc, total_budget, industry, roles, locations)

    return {
        "channel_allocations": optimized_alloc,
        "goal": goal,
        "scenario": scenario,
    }


# ═══════════════════════════════════════════════════════════════════════════════
# EXPORT FUNCTIONS
# ═══════════════════════════════════════════════════════════════════════════════

def export_scenario_excel(
    scenarios: List[Dict[str, Any]],
    comparison: Dict[str, Any],
    client_name: str = "Client",
) -> bytes:
    """Export scenarios and comparison to a styled Excel workbook.

    Sheets:
        1. Scenario Overview -- summary metrics for each scenario
        2. Per-Channel Detail -- channel-level breakdown per scenario
        3. Comparison -- side-by-side metrics with best/worst highlighting

    Design: Sapphire Blue palette, Calibri font, content at column B.

    Returns:
        Excel file bytes.
    """
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

    # Design tokens
    NAVY = "0F172A"
    SAPPHIRE = "2563EB"
    BLUE_LIGHT = "DBEAFE"
    OFF_WHITE = "F5F5F4"
    WHITE = "FFFFFF"
    STONE = "1C1917"
    GREEN = "16A34A"
    GREEN_BG = "DCFCE7"
    RED = "DC2626"
    RED_BG = "FEE2E2"
    WARM_GRAY = "E7E5E4"

    font_section = Font(name="Calibri", bold=True, size=14, color=WHITE)
    font_subsection = Font(name="Calibri", bold=True, size=12, color=NAVY)
    font_header = Font(name="Calibri", bold=True, size=10, color=WHITE)
    font_body = Font(name="Calibri", size=10, color=STONE)
    font_body_bold = Font(name="Calibri", bold=True, size=10, color=STONE)
    font_hero_value = Font(name="Calibri", bold=True, size=18, color=SAPPHIRE)
    font_hero_label = Font(name="Calibri", size=9, color="78716C")

    fill_navy = PatternFill(start_color=NAVY, end_color=NAVY, fill_type="solid")
    fill_sapphire = PatternFill(start_color=SAPPHIRE, end_color=SAPPHIRE, fill_type="solid")
    fill_light = PatternFill(start_color=BLUE_LIGHT, end_color=BLUE_LIGHT, fill_type="solid")
    fill_off_white = PatternFill(start_color=OFF_WHITE, end_color=OFF_WHITE, fill_type="solid")
    fill_white = PatternFill(start_color=WHITE, end_color=WHITE, fill_type="solid")
    fill_green = PatternFill(start_color=GREEN_BG, end_color=GREEN_BG, fill_type="solid")
    fill_red = PatternFill(start_color=RED_BG, end_color=RED_BG, fill_type="solid")

    align_center = Alignment(horizontal="center", vertical="center", wrap_text=True)
    align_left = Alignment(horizontal="left", vertical="top", wrap_text=True)
    align_right = Alignment(horizontal="right", vertical="top")

    border_thin = Border(
        left=Side(style="thin", color=WARM_GRAY),
        right=Side(style="thin", color=WARM_GRAY),
        top=Side(style="thin", color=WARM_GRAY),
        bottom=Side(style="thin", color=WARM_GRAY),
    )

    wb = Workbook()

    # ── Sheet 1: Scenario Overview ──
    ws1 = wb.active
    ws1.title = "Scenario Overview"
    ws1.sheet_properties.tabColor = SAPPHIRE

    # Column widths
    ws1.column_dimensions["A"].width = 3
    ws1.column_dimensions["B"].width = 24
    for col_letter in ["C", "D", "E"]:
        ws1.column_dimensions[col_letter].width = 22

    # Title row
    row = 2
    ws1.merge_cells(start_row=row, start_column=2, end_row=row, end_column=5)
    cell = ws1.cell(row=row, column=2, value=f"Budget Simulator -- {client_name}")
    cell.font = font_section
    cell.fill = fill_navy
    cell.alignment = align_center
    for c in range(3, 6):
        ws1.cell(row=row, column=c).fill = fill_navy

    row += 1
    ws1.merge_cells(start_row=row, start_column=2, end_row=row, end_column=5)
    cell = ws1.cell(row=row, column=2,
                    value=f"Generated: {datetime.datetime.now().strftime('%B %d, %Y')}")
    cell.font = Font(name="Calibri", italic=True, size=9, color="78716C")

    # Scenario summary cards
    row += 2
    ws1.cell(row=row, column=2, value="Metric").font = font_body_bold
    for i, sc in enumerate(scenarios):
        ws1.cell(row=row, column=3 + i, value=f"Scenario {i + 1}").font = font_header
        ws1.cell(row=row, column=3 + i).fill = fill_sapphire
        ws1.cell(row=row, column=3 + i).alignment = align_center
    ws1.cell(row=row, column=2).fill = fill_sapphire
    ws1.cell(row=row, column=2).font = font_header

    metrics = [
        ("Total Budget", "total_budget", "${:,.0f}"),
        ("Total Clicks", "total_clicks", "{:,}"),
        ("Total Applies", "total_applies", "{:,}"),
        ("Total Hires", "total_hires", "{:,}"),
        ("Blended CPC", "blended_cpc", "${:.2f}"),
        ("Blended CPA", "blended_cpa", "${:.2f}"),
        ("Blended CPH", "blended_cph", "${:,.0f}"),
        ("ROI Score", "roi_score", "{}/10"),
    ]

    for label, key, fmt in metrics:
        row += 1
        ws1.cell(row=row, column=2, value=label).font = font_body_bold
        ws1.cell(row=row, column=2).border = border_thin
        fill = fill_off_white if (row % 2 == 0) else fill_white
        ws1.cell(row=row, column=2).fill = fill

        for i, sc in enumerate(scenarios):
            val = sc.get("summary", {}).get(key, 0)
            display = fmt.format(val)
            cell = ws1.cell(row=row, column=3 + i, value=display)
            cell.font = font_body
            cell.alignment = align_center
            cell.border = border_thin
            cell.fill = fill

    # Best scenario indicator
    best_idx = comparison.get("best_scenario_index", 0)
    row += 2
    ws1.merge_cells(start_row=row, start_column=2, end_row=row, end_column=5)
    rec = comparison.get("recommendation", f"Scenario {best_idx + 1} is recommended.")
    cell = ws1.cell(row=row, column=2, value=rec)
    cell.font = Font(name="Calibri", bold=True, size=11, color=GREEN)
    cell.fill = fill_green
    cell.alignment = align_center
    for c in range(3, 6):
        ws1.cell(row=row, column=c).fill = fill_green

    # ── Sheet 2: Per-Channel Detail ──
    ws2 = wb.create_sheet("Per-Channel Detail")
    ws2.sheet_properties.tabColor = "6366F1"

    ws2.column_dimensions["A"].width = 3
    ws2.column_dimensions["B"].width = 28
    for col_letter in ["C", "D", "E", "F", "G", "H", "I"]:
        ws2.column_dimensions[col_letter].width = 16

    for sc_idx, sc in enumerate(scenarios):
        row = 2 + sc_idx * (len(CHANNEL_KEYS) + 5)

        # Scenario header
        ws2.merge_cells(start_row=row, start_column=2, end_row=row, end_column=9)
        cell = ws2.cell(row=row, column=2, value=f"Scenario {sc_idx + 1}")
        cell.font = font_section
        cell.fill = fill_navy
        for c in range(3, 10):
            ws2.cell(row=row, column=c).fill = fill_navy

        row += 1
        headers = ["Channel", "Allocation %", "Spend", "Clicks", "Applies", "Hires", "CPA", "ROI"]
        for col_idx, h in enumerate(headers):
            cell = ws2.cell(row=row, column=2 + col_idx, value=h)
            cell.font = font_header
            cell.fill = fill_sapphire
            cell.alignment = align_center

        channels = sc.get("channels", {})
        for ch_key in CHANNEL_KEYS:
            row += 1
            ch = channels.get(ch_key, {})
            fill = fill_off_white if (row % 2 == 0) else fill_white

            values = [
                ch.get("label", ch_key),
                f"{ch.get('percentage', 0):.1f}%",
                f"${ch.get('dollar_amount', 0):,.0f}",
                f"{ch.get('projected_clicks', 0):,}",
                f"{ch.get('projected_applies', 0):,}",
                f"{ch.get('projected_hires', 0):,}",
                f"${ch.get('cpa', 0):,.2f}",
                f"{ch.get('roi_score', 0)}/10",
            ]
            for col_idx, v in enumerate(values):
                cell = ws2.cell(row=row, column=2 + col_idx, value=v)
                cell.font = font_body
                cell.alignment = align_center if col_idx > 0 else align_left
                cell.border = border_thin
                cell.fill = fill

    # ── Sheet 3: Comparison ──
    ws3 = wb.create_sheet("Comparison")
    ws3.sheet_properties.tabColor = "16A34A"

    ws3.column_dimensions["A"].width = 3
    ws3.column_dimensions["B"].width = 24
    for col_letter in ["C", "D", "E", "F"]:
        ws3.column_dimensions[col_letter].width = 20

    row = 2
    ws3.merge_cells(start_row=row, start_column=2, end_row=row, end_column=6)
    cell = ws3.cell(row=row, column=2, value="Scenario Comparison")
    cell.font = font_section
    cell.fill = fill_navy
    for c in range(3, 7):
        ws3.cell(row=row, column=c).fill = fill_navy

    row += 2
    ws3.cell(row=row, column=2, value="Metric").font = font_header
    ws3.cell(row=row, column=2).fill = fill_sapphire
    for i in range(len(scenarios)):
        cell = ws3.cell(row=row, column=3 + i, value=f"Scenario {i + 1}")
        cell.font = font_header
        cell.fill = fill_sapphire
        cell.alignment = align_center
    ws3.cell(row=row, column=3 + len(scenarios), value="Best").font = font_header
    ws3.cell(row=row, column=3 + len(scenarios)).fill = fill_sapphire
    ws3.cell(row=row, column=3 + len(scenarios)).alignment = align_center

    mbw = comparison.get("metrics_best_worst", {})
    for label, key, fmt in metrics:
        row += 1
        ws3.cell(row=row, column=2, value=label).font = font_body_bold
        ws3.cell(row=row, column=2).border = border_thin

        best_for_metric = mbw.get(key, {}).get("best", 0)
        worst_for_metric = mbw.get(key, {}).get("worst", -1)

        for i, sc in enumerate(scenarios):
            val = sc.get("summary", {}).get(key, 0)
            display = fmt.format(val)
            cell = ws3.cell(row=row, column=3 + i, value=display)
            cell.font = font_body
            cell.alignment = align_center
            cell.border = border_thin

            if i == best_for_metric and len(scenarios) > 1:
                cell.fill = fill_green
                cell.font = Font(name="Calibri", bold=True, size=10, color=GREEN)
            elif i == worst_for_metric and len(scenarios) > 1:
                cell.fill = fill_red

        # Best indicator
        cell = ws3.cell(row=row, column=3 + len(scenarios),
                        value=f"Scenario {best_for_metric + 1}")
        cell.font = Font(name="Calibri", bold=True, size=10, color=GREEN)
        cell.alignment = align_center
        cell.border = border_thin

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def export_scenario_ppt(
    scenarios: List[Dict[str, Any]],
    comparison: Dict[str, Any],
    client_name: str = "Client",
) -> bytes:
    """Export scenarios and comparison to a Nova AI Suite branded PowerPoint.

    Slides:
        1. Title slide
        2-4. Scenario 1/2/3 breakdown (one per scenario)
        5. Comparison slide
        6. Recommendation slide

    Returns:
        PPT file bytes.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    # Brand colors
    NAVY = RGBColor(0x20, 0x20, 0x58)
    BLUE = RGBColor(0x5A, 0x54, 0xBD)
    TEAL = RGBColor(0x6B, 0xB3, 0xCD)
    WHITE_C = RGBColor(0xFF, 0xFF, 0xFF)
    OFF_WHITE_C = RGBColor(0xFF, 0xFD, 0xF9)
    DARK_TEXT = RGBColor(0x20, 0x20, 0x58)
    MUTED = RGBColor(0x59, 0x67, 0x80)
    GREEN_C = RGBColor(0x33, 0x87, 0x21)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def _add_bg(slide, color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_text_box(slide, left, top, width, height, text, font_size=12,
                      bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT):
        from pptx.enum.shapes import MSO_SHAPE
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                         Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = alignment
        return txBox

    # ── Slide 1: Title ──
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_bg(slide1, NAVY)
    _add_text_box(slide1, 1, 1.5, 11, 1, "BUDGET SIMULATOR",
                  font_size=36, bold=True, color=WHITE_C, alignment=PP_ALIGN.CENTER)
    _add_text_box(slide1, 1, 3, 11, 0.8, f"Multi-Scenario Analysis for {client_name}",
                  font_size=20, color=TEAL, alignment=PP_ALIGN.CENTER)
    _add_text_box(slide1, 1, 4.5, 11, 0.5,
                  f"Generated {datetime.datetime.now().strftime('%B %d, %Y')} | Powered by Nova AI Suite",
                  font_size=12, color=MUTED, alignment=PP_ALIGN.CENTER)

    # ── Slides 2-4: Scenario breakdown ──
    for sc_idx, sc in enumerate(scenarios):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_bg(slide, OFF_WHITE_C)

        _add_text_box(slide, 0.5, 0.3, 12, 0.6,
                      f"Scenario {sc_idx + 1}",
                      font_size=24, bold=True, color=NAVY)

        summary = sc.get("summary", {})
        # Hero metrics
        hero_metrics = [
            ("Hires", f"{summary.get('total_hires', 0):,}"),
            ("Applies", f"{summary.get('total_applies', 0):,}"),
            ("CPA", f"${summary.get('blended_cpa', 0):,.2f}"),
            ("ROI", f"{summary.get('roi_score', 0)}/10"),
        ]

        x_start = 0.5
        for mi, (mlabel, mval) in enumerate(hero_metrics):
            x = x_start + mi * 3.1
            # Metric card background
            from pptx.enum.shapes import MSO_SHAPE
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(1.1), Inches(2.8), Inches(1.2)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = WHITE_C
            shape.line.color.rgb = RGBColor(0xDD, 0xDB, 0xFF)
            shape.line.width = Pt(1)

            _add_text_box(slide, x + 0.2, 1.2, 2.4, 0.3, mlabel,
                          font_size=10, color=MUTED)
            _add_text_box(slide, x + 0.2, 1.5, 2.4, 0.5, mval,
                          font_size=22, bold=True, color=BLUE)

        # Channel table
        channels = sc.get("channels", {})
        _add_text_box(slide, 0.5, 2.6, 12, 0.4, "Channel Allocation",
                      font_size=14, bold=True, color=NAVY)

        y = 3.1
        # Header row
        headers = ["Channel", "%", "Spend", "Clicks", "Applies", "Hires", "CPA"]
        col_widths = [3.0, 1.0, 1.5, 1.5, 1.5, 1.2, 1.5]
        x = 0.5
        for hi, h in enumerate(headers):
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(y), Inches(col_widths[hi]), Inches(0.35)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = BLUE
            shape.line.fill.background()
            tf = shape.text_frame
            p = tf.paragraphs[0]
            p.text = h
            p.font.size = Pt(8)
            p.font.bold = True
            p.font.color.rgb = WHITE_C
            p.alignment = PP_ALIGN.CENTER
            x += col_widths[hi]

        y += 0.35
        for ch_key in CHANNEL_KEYS:
            ch = channels.get(ch_key, {})
            row_data = [
                ch.get("label", ch_key),
                f"{ch.get('percentage', 0):.0f}%",
                f"${ch.get('dollar_amount', 0):,.0f}",
                f"{ch.get('projected_clicks', 0):,}",
                f"{ch.get('projected_applies', 0):,}",
                f"{ch.get('projected_hires', 0):,}",
                f"${ch.get('cpa', 0):,.0f}",
            ]

            x = 0.5
            for ci, val in enumerate(row_data):
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(x), Inches(y), Inches(col_widths[ci]), Inches(0.3)
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = OFF_WHITE_C if (CHANNEL_KEYS.index(ch_key) % 2 == 0) else WHITE_C
                shape.line.fill.background()
                tf = shape.text_frame
                p = tf.paragraphs[0]
                p.text = val
                p.font.size = Pt(7)
                p.font.color.rgb = DARK_TEXT
                p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
                x += col_widths[ci]
            y += 0.3

    # ── Slide 5: Comparison ──
    slide_comp = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide_comp, OFF_WHITE_C)
    _add_text_box(slide_comp, 0.5, 0.3, 12, 0.6,
                  "Scenario Comparison", font_size=24, bold=True, color=NAVY)

    comp_scenarios = comparison.get("scenarios", [])
    comp_metrics = [
        ("Total Hires", "total_hires"),
        ("Total Applies", "total_applies"),
        ("Blended CPA", "blended_cpa"),
        ("Blended CPC", "blended_cpc"),
        ("ROI Score", "roi_score"),
    ]

    y = 1.2
    # Header
    x = 0.5
    for header in ["Metric"] + [f"Scenario {i+1}" for i in range(len(comp_scenarios))]:
        w = 3.5 if header == "Metric" else 2.5
        shape = slide_comp.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.4)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = BLUE
        shape.line.fill.background()
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = header
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = WHITE_C
        p.alignment = PP_ALIGN.CENTER
        x += w

    y += 0.4
    best_idx = comparison.get("best_scenario_index", 0)
    mbw = comparison.get("metrics_best_worst", {})

    for mlabel, mkey in comp_metrics:
        x = 0.5
        # Metric label
        shape = slide_comp.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(3.5), Inches(0.4)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE_C
        shape.line.fill.background()
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = mlabel
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = DARK_TEXT
        x += 3.5

        best_for = mbw.get(mkey, {}).get("best", 0)
        for si, s in enumerate(comp_scenarios):
            val = s.get(mkey, 0)
            if mkey in ("blended_cpa", "blended_cpc", "blended_cph"):
                display = f"${val:,.2f}"
            elif mkey == "roi_score":
                display = f"{val}/10"
            else:
                display = f"{val:,}"

            shape = slide_comp.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(2.5), Inches(0.4)
            )
            shape.fill.solid()
            is_best = (si == best_for)
            shape.fill.fore_color.rgb = RGBColor(0xE6, 0xF2, 0xE0) if is_best else WHITE_C
            shape.line.fill.background()
            tf = shape.text_frame
            p = tf.paragraphs[0]
            p.text = display
            p.font.size = Pt(10)
            p.font.bold = is_best
            p.font.color.rgb = GREEN_C if is_best else DARK_TEXT
            p.alignment = PP_ALIGN.CENTER
            x += 2.5
        y += 0.4

    # ── Slide 6: Recommendation ──
    slide_rec = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide_rec, NAVY)
    _add_text_box(slide_rec, 1, 1.5, 11, 1, "RECOMMENDATION",
                  font_size=32, bold=True, color=WHITE_C, alignment=PP_ALIGN.CENTER)

    rec_text = comparison.get("recommendation", f"Scenario {best_idx + 1} is the recommended allocation.")
    _add_text_box(slide_rec, 1.5, 3, 10, 1.5, rec_text,
                  font_size=16, color=TEAL, alignment=PP_ALIGN.CENTER)

    if comp_scenarios and best_idx < len(comp_scenarios):
        best_sc = comp_scenarios[best_idx]
        details = (
            f"Projected Hires: {best_sc.get('total_hires', 0):,}  |  "
            f"CPA: ${best_sc.get('blended_cpa', 0):,.2f}  |  "
            f"ROI: {best_sc.get('roi_score', 0)}/10"
        )
        _add_text_box(slide_rec, 1.5, 4.8, 10, 0.6, details,
                      font_size=14, bold=True, color=WHITE_C, alignment=PP_ALIGN.CENTER)

    _add_text_box(slide_rec, 1, 6.2, 11, 0.4,
                  "Powered by Nova AI Suite | Budget Simulator",
                  font_size=10, color=MUTED, alignment=PP_ALIGN.CENTER)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
