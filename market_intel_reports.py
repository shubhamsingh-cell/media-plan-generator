"""
market_intel_reports.py -- Market Intelligence Reports for Recruitment Advertising

Generates comprehensive market research reports for specific industries and roles,
combining labor market data, compensation benchmarks, channel performance metrics,
CPC trends, seasonal patterns, competitor analysis, talent supply data, and
actionable recommendations.

Outputs:
    - Structured Python dict (JSON-serializable)
    - Excel workbook (Sapphire Blue palette: #0F172A, #2563EB, #DBEAFE)
    - PowerPoint deck (Joveo branding: Port Gore #202058, Blue Violet #5A54BD, Downy #6BB3CD)

Entry points:
    generate_report(data)              -- Main report generator
    generate_intel_excel(report_data)  -- Excel export (returns bytes)
    generate_intel_ppt(report_data)    -- PPT export (returns bytes)
    handle_market_intel_request(path, method, body)  -- Unified HTTP handler

All external module imports are lazy with fallback data so the module
degrades gracefully when dependencies are unavailable.
"""

from __future__ import annotations

import io
import json
import logging
import math
import statistics
import threading
from concurrent.futures import ThreadPoolExecutor, as_completed
from datetime import datetime
from typing import Any, Dict, List, Optional, Tuple

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Lazy imports with thread-safe loading
# ---------------------------------------------------------------------------

_load_lock = threading.Lock()
_loaded = False
_trend_engine = None
_research = None
_api_enrichment = None
_collar_intelligence = None
_shared_utils = None
_INDUSTRY_LABEL_MAP: Dict[str, str] = {}


def _lazy_load():
    """Import optional modules once, thread-safe."""
    global _loaded, _trend_engine, _research, _api_enrichment
    global _collar_intelligence, _shared_utils, _INDUSTRY_LABEL_MAP

    if _loaded:
        return
    with _load_lock:
        if _loaded:
            return
        try:
            import trend_engine as _te
            _trend_engine = _te
        except Exception:
            pass
        try:
            import research as _r
            _research = _r
        except Exception:
            pass
        try:
            import api_enrichment as _ae
            _api_enrichment = _ae
        except Exception:
            pass
        try:
            import collar_intelligence as _ci
            _collar_intelligence = _ci
        except Exception:
            pass
        try:
            import shared_utils as _su
            _shared_utils = _su
            _INDUSTRY_LABEL_MAP = getattr(_su, "INDUSTRY_LABEL_MAP", {})
        except Exception:
            pass
        _loaded = True


# ═══════════════════════════════════════════════════════════════════════════════
# KNOWLEDGE BASE -- 22 Industries
# ═══════════════════════════════════════════════════════════════════════════════

INDUSTRY_KB: Dict[str, Dict[str, Any]] = {
    "technology": {
        "label": "Technology & Software", "collar": "white_collar",
        "growth_rate": 0.15, "avg_time_to_fill": 42, "turnover_rate": 0.132,
        "top_roles": ["Software Engineer", "Data Scientist", "Product Manager", "DevOps Engineer"],
        "talent_supply": "moderate", "demand_trend": "rising",
    },
    "healthcare_medical": {
        "label": "Healthcare & Medical", "collar": "grey_collar",
        "growth_rate": 0.13, "avg_time_to_fill": 49, "turnover_rate": 0.195,
        "top_roles": ["Registered Nurse", "Medical Assistant", "Physician", "Pharmacist"],
        "talent_supply": "tight", "demand_trend": "rising",
    },
    "manufacturing": {
        "label": "Manufacturing", "collar": "blue_collar",
        "growth_rate": 0.03, "avg_time_to_fill": 33, "turnover_rate": 0.28,
        "top_roles": ["Machine Operator", "Quality Inspector", "Welder", "CNC Machinist"],
        "talent_supply": "tight", "demand_trend": "stable",
    },
    "finance_banking": {
        "label": "Finance & Banking", "collar": "white_collar",
        "growth_rate": 0.07, "avg_time_to_fill": 44, "turnover_rate": 0.11,
        "top_roles": ["Financial Analyst", "Accountant", "Compliance Officer", "Risk Manager"],
        "talent_supply": "adequate", "demand_trend": "stable",
    },
    "retail_consumer": {
        "label": "Retail & Consumer", "collar": "pink_collar",
        "growth_rate": 0.02, "avg_time_to_fill": 21, "turnover_rate": 0.60,
        "top_roles": ["Store Manager", "Sales Associate", "Merchandiser", "Cashier"],
        "talent_supply": "surplus", "demand_trend": "stable",
    },
    "logistics_supply_chain": {
        "label": "Logistics & Supply Chain", "collar": "blue_collar",
        "growth_rate": 0.08, "avg_time_to_fill": 25, "turnover_rate": 0.43,
        "top_roles": ["Warehouse Associate", "Truck Driver", "Forklift Operator", "Logistics Coordinator"],
        "talent_supply": "tight", "demand_trend": "rising",
    },
    "construction_real_estate": {
        "label": "Construction & Real Estate", "collar": "blue_collar",
        "growth_rate": 0.04, "avg_time_to_fill": 30, "turnover_rate": 0.35,
        "top_roles": ["Electrician", "Plumber", "Project Manager", "Carpenter"],
        "talent_supply": "tight", "demand_trend": "stable",
    },
    "hospitality_travel": {
        "label": "Hospitality & Travel", "collar": "pink_collar",
        "growth_rate": 0.10, "avg_time_to_fill": 18, "turnover_rate": 0.73,
        "top_roles": ["Hotel Manager", "Chef", "Front Desk Agent", "Server"],
        "talent_supply": "moderate", "demand_trend": "rising",
    },
    "education": {
        "label": "Education", "collar": "white_collar",
        "growth_rate": 0.05, "avg_time_to_fill": 55, "turnover_rate": 0.16,
        "top_roles": ["Teacher", "Professor", "School Administrator", "Tutor"],
        "talent_supply": "moderate", "demand_trend": "stable",
    },
    "pharma_biotech": {
        "label": "Pharma & Biotech", "collar": "white_collar",
        "growth_rate": 0.11, "avg_time_to_fill": 52, "turnover_rate": 0.10,
        "top_roles": ["Research Scientist", "Clinical Research Associate", "Regulatory Affairs", "QA Analyst"],
        "talent_supply": "tight", "demand_trend": "rising",
    },
    "energy_utilities": {
        "label": "Energy & Utilities", "collar": "blue_collar",
        "growth_rate": 0.06, "avg_time_to_fill": 38, "turnover_rate": 0.12,
        "top_roles": ["Electrical Engineer", "Lineworker", "Plant Operator", "Safety Inspector"],
        "talent_supply": "moderate", "demand_trend": "stable",
    },
    "automotive": {
        "label": "Automotive", "collar": "blue_collar",
        "growth_rate": 0.04, "avg_time_to_fill": 32, "turnover_rate": 0.22,
        "top_roles": ["Auto Technician", "Assembly Worker", "Design Engineer", "Parts Advisor"],
        "talent_supply": "moderate", "demand_trend": "stable",
    },
    "insurance": {
        "label": "Insurance", "collar": "white_collar",
        "growth_rate": 0.05, "avg_time_to_fill": 40, "turnover_rate": 0.12,
        "top_roles": ["Claims Adjuster", "Underwriter", "Actuary", "Insurance Agent"],
        "talent_supply": "adequate", "demand_trend": "stable",
    },
    "aerospace_defense": {
        "label": "Aerospace & Defense", "collar": "white_collar",
        "growth_rate": 0.06, "avg_time_to_fill": 58, "turnover_rate": 0.09,
        "top_roles": ["Aerospace Engineer", "Systems Analyst", "Avionics Tech", "Program Manager"],
        "talent_supply": "tight", "demand_trend": "rising",
    },
    "food_beverage": {
        "label": "Food & Beverage", "collar": "blue_collar",
        "growth_rate": 0.03, "avg_time_to_fill": 19, "turnover_rate": 0.82,
        "top_roles": ["Line Cook", "Food Production Worker", "Quality Control", "Delivery Driver"],
        "talent_supply": "surplus", "demand_trend": "stable",
    },
    "telecommunications": {
        "label": "Telecommunications", "collar": "white_collar",
        "growth_rate": 0.04, "avg_time_to_fill": 36, "turnover_rate": 0.15,
        "top_roles": ["Network Engineer", "Field Technician", "Sales Executive", "RF Engineer"],
        "talent_supply": "moderate", "demand_trend": "stable",
    },
    "legal_services": {
        "label": "Legal Services", "collar": "white_collar",
        "growth_rate": 0.06, "avg_time_to_fill": 48, "turnover_rate": 0.17,
        "top_roles": ["Paralegal", "Associate Attorney", "Legal Assistant", "Compliance Analyst"],
        "talent_supply": "adequate", "demand_trend": "stable",
    },
    "media_entertainment": {
        "label": "Media & Entertainment", "collar": "white_collar",
        "growth_rate": 0.08, "avg_time_to_fill": 35, "turnover_rate": 0.20,
        "top_roles": ["Content Creator", "Video Editor", "Marketing Manager", "UX Designer"],
        "talent_supply": "surplus", "demand_trend": "rising",
    },
    "mental_health": {
        "label": "Mental Health & Behavioral", "collar": "grey_collar",
        "growth_rate": 0.18, "avg_time_to_fill": 55, "turnover_rate": 0.25,
        "top_roles": ["Therapist", "Counselor", "Psychiatrist", "Social Worker"],
        "talent_supply": "tight", "demand_trend": "rising",
    },
    "government": {
        "label": "Government & Public Sector", "collar": "white_collar",
        "growth_rate": 0.02, "avg_time_to_fill": 65, "turnover_rate": 0.08,
        "top_roles": ["Program Analyst", "Administrative Officer", "IT Specialist", "Budget Analyst"],
        "talent_supply": "adequate", "demand_trend": "stable",
    },
    "nonprofit": {
        "label": "Nonprofit & Social Services", "collar": "pink_collar",
        "growth_rate": 0.04, "avg_time_to_fill": 40, "turnover_rate": 0.19,
        "top_roles": ["Program Manager", "Grant Writer", "Case Manager", "Development Director"],
        "talent_supply": "adequate", "demand_trend": "stable",
    },
    "agriculture": {
        "label": "Agriculture & Farming", "collar": "blue_collar",
        "growth_rate": 0.01, "avg_time_to_fill": 22, "turnover_rate": 0.45,
        "top_roles": ["Farm Worker", "Agricultural Technician", "Irrigation Specialist", "Equipment Operator"],
        "talent_supply": "tight", "demand_trend": "declining",
    },
}

# ═══════════════════════════════════════════════════════════════════════════════
# REGIONAL SALARY DATA -- 35 Metros
# ═══════════════════════════════════════════════════════════════════════════════

METRO_SALARY_INDEX: Dict[str, Dict[str, Any]] = {
    "new york": {"coli": 130, "median": 72000, "p25": 52000, "p75": 105000, "p90": 145000},
    "san francisco": {"coli": 145, "median": 85000, "p25": 60000, "p75": 125000, "p90": 170000},
    "los angeles": {"coli": 120, "median": 65000, "p25": 46000, "p75": 95000, "p90": 135000},
    "chicago": {"coli": 107, "median": 60000, "p25": 42000, "p75": 85000, "p90": 120000},
    "houston": {"coli": 96, "median": 58000, "p25": 40000, "p75": 82000, "p90": 115000},
    "dallas": {"coli": 98, "median": 59000, "p25": 41000, "p75": 84000, "p90": 118000},
    "austin": {"coli": 105, "median": 65000, "p25": 45000, "p75": 92000, "p90": 130000},
    "seattle": {"coli": 130, "median": 78000, "p25": 55000, "p75": 115000, "p90": 155000},
    "boston": {"coli": 125, "median": 72000, "p25": 50000, "p75": 100000, "p90": 140000},
    "washington dc": {"coli": 122, "median": 75000, "p25": 52000, "p75": 105000, "p90": 142000},
    "denver": {"coli": 108, "median": 63000, "p25": 44000, "p75": 90000, "p90": 125000},
    "atlanta": {"coli": 100, "median": 58000, "p25": 40000, "p75": 82000, "p90": 115000},
    "miami": {"coli": 112, "median": 55000, "p25": 38000, "p75": 78000, "p90": 110000},
    "phoenix": {"coli": 100, "median": 55000, "p25": 38000, "p75": 78000, "p90": 108000},
    "philadelphia": {"coli": 110, "median": 62000, "p25": 43000, "p75": 88000, "p90": 122000},
    "san diego": {"coli": 118, "median": 65000, "p25": 45000, "p75": 92000, "p90": 128000},
    "minneapolis": {"coli": 105, "median": 62000, "p25": 43000, "p75": 88000, "p90": 120000},
    "detroit": {"coli": 92, "median": 52000, "p25": 36000, "p75": 74000, "p90": 102000},
    "portland": {"coli": 110, "median": 62000, "p25": 43000, "p75": 88000, "p90": 120000},
    "nashville": {"coli": 100, "median": 56000, "p25": 39000, "p75": 80000, "p90": 112000},
    "charlotte": {"coli": 98, "median": 57000, "p25": 40000, "p75": 80000, "p90": 112000},
    "raleigh": {"coli": 100, "median": 60000, "p25": 42000, "p75": 85000, "p90": 118000},
    "salt lake city": {"coli": 100, "median": 56000, "p25": 39000, "p75": 80000, "p90": 110000},
    "columbus": {"coli": 92, "median": 54000, "p25": 38000, "p75": 76000, "p90": 105000},
    "indianapolis": {"coli": 90, "median": 52000, "p25": 36000, "p75": 74000, "p90": 102000},
    "san antonio": {"coli": 90, "median": 50000, "p25": 35000, "p75": 72000, "p90": 100000},
    "london": {"coli": 135, "median": 48000, "p25": 34000, "p75": 70000, "p90": 98000},
    "toronto": {"coli": 108, "median": 55000, "p25": 38000, "p75": 78000, "p90": 108000},
    "sydney": {"coli": 115, "median": 62000, "p25": 43000, "p75": 88000, "p90": 120000},
    "berlin": {"coli": 95, "median": 50000, "p25": 35000, "p75": 72000, "p90": 98000},
    "singapore": {"coli": 120, "median": 55000, "p25": 38000, "p75": 78000, "p90": 110000},
    "mumbai": {"coli": 30, "median": 12000, "p25": 7000, "p75": 22000, "p90": 38000},
    "bangalore": {"coli": 28, "median": 14000, "p25": 8000, "p75": 25000, "p90": 42000},
    "dubai": {"coli": 85, "median": 40000, "p25": 28000, "p75": 60000, "p90": 85000},
    "tokyo": {"coli": 110, "median": 42000, "p25": 30000, "p75": 62000, "p90": 88000},
}

# ═══════════════════════════════════════════════════════════════════════════════
# PLATFORM PERFORMANCE BENCHMARKS
# ═══════════════════════════════════════════════════════════════════════════════

PLATFORM_BENCHMARKS: Dict[str, Dict[str, float]] = {
    "indeed":        {"cpc": 0.45, "cpa": 18.0, "ctr": 0.032, "conv_rate": 0.025},
    "linkedin":      {"cpc": 5.50, "cpa": 85.0, "ctr": 0.028, "conv_rate": 0.065},
    "google_search": {"cpc": 3.80, "cpa": 55.0, "ctr": 0.042, "conv_rate": 0.069},
    "meta":          {"cpc": 0.90, "cpa": 32.0, "ctr": 0.015, "conv_rate": 0.047},
    "programmatic":  {"cpc": 0.65, "cpa": 22.0, "ctr": 0.038, "conv_rate": 0.029},
    "ziprecruiter":  {"cpc": 0.55, "cpa": 20.0, "ctr": 0.035, "conv_rate": 0.028},
    "glassdoor":     {"cpc": 1.20, "cpa": 38.0, "ctr": 0.030, "conv_rate": 0.032},
    "tiktok":        {"cpc": 0.70, "cpa": 28.0, "ctr": 0.012, "conv_rate": 0.018},
}

# Collar-type CPC multipliers
_COLLAR_CPC_MULT: Dict[str, float] = {
    "blue_collar": 0.65, "white_collar": 1.25,
    "grey_collar": 1.10, "pink_collar": 0.80,
}

_SEASONAL_MULT: Dict[int, float] = {
    1: 1.12, 2: 1.08, 3: 1.05, 4: 1.02, 5: 0.98, 6: 0.95,
    7: 0.92, 8: 0.94, 9: 1.06, 10: 1.03, 11: 0.96, 12: 0.89,
}

_FALLBACK_SALARY: Dict[str, Any] = {
    "p10": 35000, "p25": 45000, "median": 58000, "p75": 78000, "p90": 105000,
}


# ═══════════════════════════════════════════════════════════════════════════════
# UTILITY HELPERS
# ═══════════════════════════════════════════════════════════════════════════════

def _industry_label(key: str) -> str:
    _lazy_load()
    m = _INDUSTRY_LABEL_MAP or {}
    if key in m:
        return m[key]
    kb = INDUSTRY_KB.get(key)
    if kb:
        return kb["label"]
    return key.replace("_", " ").title()


def _detect_collar(role: str, industry: str) -> str:
    """Classify into blue/white/grey/pink collar."""
    _lazy_load()

    # Keyword heuristics first (most reliable for role-level detection)
    blue_kw = {"driver", "warehouse", "technician", "mechanic", "operator",
               "welder", "electrician", "plumber", "carpenter", "assembler",
               "laborer", "forklift", "janitor", "maintenance", "hvac",
               "construction", "installer", "loader"}
    grey_kw = {"nurse", "therapist", "paramedic", "emt", "dental", "radiology",
               "phlebotom", "respiratory", "surgical tech"}
    pink_kw = {"receptionist", "secretary", "cashier", "caregiver", "nanny",
               "aide", "clerk", "librarian"}
    rl = role.lower()
    for kw in grey_kw:
        if kw in rl:
            return "grey_collar"
    for kw in blue_kw:
        if kw in rl:
            return "blue_collar"
    for kw in pink_kw:
        if kw in rl:
            return "pink_collar"
    # Try collar_intelligence module
    if _collar_intelligence:
        try:
            result = _collar_intelligence.classify_collar(role)
            ci_collar = result.get("collar", "")
            if ci_collar in ("blue_collar", "white_collar", "grey_collar", "pink_collar"):
                return ci_collar
        except Exception:
            pass
    # Fallback to industry KB
    kb = INDUSTRY_KB.get(industry)
    if kb:
        return kb.get("collar", "white_collar")
    return "white_collar"


def _safe_get(func, *args, default=None, **kwargs):
    try:
        r = func(*args, **kwargs)
        return r if r is not None else default
    except Exception as exc:
        logger.debug("_safe_get: %s: %s", type(exc).__name__, exc)
        return default


def _fmt_currency(v: float, d: int = 0) -> str:
    if d == 0:
        return f"${v:,.0f}"
    return f"${v:,.{d}f}"


def _fmt_pct(v: float, d: int = 1) -> str:
    return f"{v * 100:.{d}f}%" if v < 1 else f"{v:.{d}f}%"


def _resolve_metro(locations: List[str]) -> Optional[Dict[str, Any]]:
    """Find the best matching metro salary data for a list of locations."""
    for loc in locations:
        key = loc.lower().split(",")[0].strip()
        if key in METRO_SALARY_INDEX:
            return METRO_SALARY_INDEX[key]
        for metro_key in METRO_SALARY_INDEX:
            if metro_key in key or key in metro_key:
                return METRO_SALARY_INDEX[metro_key]
    return None


# ═══════════════════════════════════════════════════════════════════════════════
# DATA COLLECTION
# ═══════════════════════════════════════════════════════════════════════════════

def _collect_market_overview(industry: str, locations: List[str]) -> Dict[str, Any]:
    """Gather industry hiring trends, growth, talent supply/demand."""
    _lazy_load()
    kb = INDUSTRY_KB.get(industry, INDUSTRY_KB.get("technology", {}))
    data: Dict[str, Any] = {
        "industry": industry,
        "industry_label": _industry_label(industry),
        "growth_rate": kb.get("growth_rate", 0.05),
        "avg_time_to_fill": kb.get("avg_time_to_fill", 40),
        "turnover_rate": kb.get("turnover_rate", 0.15),
        "top_roles": kb.get("top_roles", []),
        "talent_supply": kb.get("talent_supply", "moderate"),
        "demand_trend": kb.get("demand_trend", "stable"),
        "source": "knowledge_base",
    }
    if _research:
        lmi = _safe_get(_research.get_labour_market_intelligence, industry, locations, default={})
        if lmi:
            data["research_lmi"] = lmi
            data["source"] = "research+knowledge_base"
    if _api_enrichment:
        jolts = _safe_get(_api_enrichment.get_jolts_hiring_difficulty, industry, default={})
        if jolts:
            data["jolts"] = jolts
            data["source"] = "api+" + data["source"]
    # Difficulty score
    score = 50
    supply = data["talent_supply"]
    if supply == "tight":
        score += 20
    elif supply == "surplus":
        score -= 15
    if data["demand_trend"] == "rising":
        score += 10
    elif data["demand_trend"] == "declining":
        score -= 10
    data["difficulty_score"] = max(0, min(100, score))
    return data


def _collect_salary_benchmarks(
    role_category: str, industry: str, locations: List[str],
) -> Dict[str, Any]:
    """Salary ranges by location and percentile."""
    _lazy_load()
    data: Dict[str, Any] = {"source": "knowledge_base", "by_location": {}}

    metro = _resolve_metro(locations)
    base = metro if metro else dict(_FALLBACK_SALARY)

    # Apply industry multiplier
    kb = INDUSTRY_KB.get(industry, {})
    collar = kb.get("collar", "white_collar")
    mult = {"white_collar": 1.15, "grey_collar": 1.0, "blue_collar": 0.85, "pink_collar": 0.80}.get(collar, 1.0)

    for loc in locations:
        loc_metro = _resolve_metro([loc]) or base
        loc_data = {
            "p10": round(loc_metro.get("p10", base.get("p10", 35000)) * mult),
            "p25": round(loc_metro.get("p25", base.get("p25", 45000)) * mult),
            "median": round(loc_metro.get("median", base.get("median", 58000)) * mult),
            "p75": round(loc_metro.get("p75", base.get("p75", 78000)) * mult),
            "p90": round(loc_metro.get("p90", base.get("p90", 105000)) * mult),
            "coli": loc_metro.get("coli", 100),
        }
        data["by_location"][loc] = loc_data

    # BLS/API enrichment
    if _api_enrichment:
        bls = _safe_get(_api_enrichment.fetch_salary_data, [role_category], default={})
        if bls:
            data["bls_salary"] = bls
            data["source"] = "api+knowledge_base"

    # Aggregate
    all_medians = [v["median"] for v in data["by_location"].values() if v.get("median")]
    data["aggregate"] = {
        "median": round(statistics.mean(all_medians)) if all_medians else base.get("median", 58000),
        "min_p10": min((v["p10"] for v in data["by_location"].values()), default=35000),
        "max_p90": max((v["p90"] for v in data["by_location"].values()), default=105000),
    }
    return data


def _collect_channel_performance(industry: str, collar_type: str) -> Dict[str, Any]:
    """Platform effectiveness data with CPC, CPA, conversion rates."""
    _lazy_load()
    data: Dict[str, Any] = {"source": "knowledge_base", "channels": {}}

    collar_mult = _COLLAR_CPC_MULT.get(collar_type, 1.0)

    if _trend_engine:
        benchmarks = _safe_get(_trend_engine.get_all_platform_benchmarks, industry, collar_type, default={})
        if benchmarks:
            data["source"] = "trend_engine"
            for plat, m in benchmarks.items():
                data["channels"][plat] = {
                    "cpc": m.get("cpc", {}).get("value", 0),
                    "cpa": m.get("cpa", {}).get("value", 0) if "cpa" in m else None,
                    "ctr": m.get("ctr", {}).get("value", 0) if "ctr" in m else None,
                    "conv_rate": m.get("conv_rate", {}).get("value", 0) if "conv_rate" in m else None,
                    "confidence": m.get("cpc", {}).get("confidence", "medium"),
                    "trend": m.get("cpc", {}).get("trend_direction", "stable"),
                }

    if not data["channels"]:
        data["source"] = "knowledge_base"
        for plat, bm in PLATFORM_BENCHMARKS.items():
            data["channels"][plat] = {
                "cpc": round(bm["cpc"] * collar_mult, 2),
                "cpa": round(bm["cpa"] * collar_mult, 2),
                "ctr": bm.get("ctr"),
                "conv_rate": bm.get("conv_rate"),
                "confidence": "medium",
                "trend": "stable",
            }

    data["ranked_by_cpc"] = sorted(data["channels"].items(), key=lambda x: x[1].get("cpc", 999))
    data["ranked_by_cpa"] = sorted(data["channels"].items(), key=lambda x: x[1].get("cpa") or 999)
    return data


def _collect_competitor_analysis(competitors: List[str], industry: str) -> Dict[str, Any]:
    """Competitor hiring activity and strategies."""
    _lazy_load()
    data: Dict[str, Any] = {"source": "estimate", "competitors": []}
    kb = INDUSTRY_KB.get(industry, {})
    top_roles = kb.get("top_roles", ["General"])

    for comp_name in competitors[:10]:
        entry: Dict[str, Any] = {
            "name": comp_name,
            "estimated_openings": 50 + hash(comp_name) % 200,
            "primary_channels": ["Indeed", "LinkedIn", "Glassdoor"],
            "top_roles_hiring": top_roles[:3],
            "employer_brand_strength": "medium",
        }
        # Try Wikipedia enrichment
        if _api_enrichment and hasattr(_api_enrichment, "fetch_company_info"):
            info = _safe_get(_api_enrichment.fetch_company_info, comp_name, default={})
            if info and info.get("description"):
                entry["description"] = info["description"][:300]
                data["source"] = "api+estimate"
        data["competitors"].append(entry)

    data["market_position_summary"] = (
        f"Analysis covers {len(data['competitors'])} competitor(s) in the "
        f"{_industry_label(industry)} sector."
    )
    return data


def _collect_talent_supply(role_category: str, industry: str, locations: List[str]) -> Dict[str, Any]:
    """Talent pool size, availability, mobility patterns."""
    kb = INDUSTRY_KB.get(industry, {})
    supply_level = kb.get("talent_supply", "moderate")
    pool_mult = {"tight": 0.6, "moderate": 1.0, "adequate": 1.3, "surplus": 1.8}.get(supply_level, 1.0)
    base_pool = 5000

    data: Dict[str, Any] = {
        "source": "estimate",
        "supply_level": supply_level,
        "estimated_pool_size": int(base_pool * pool_mult * len(locations)),
        "active_seekers_pct": 0.25 if supply_level == "tight" else 0.40,
        "passive_candidates_pct": 0.55,
        "mobility_index": 0.7 if supply_level == "tight" else 0.5,
        "avg_tenure_years": 3.2 if kb.get("turnover_rate", 0.15) > 0.30 else 5.5,
        "remote_eligible_pct": 0.65 if kb.get("collar", "white_collar") == "white_collar" else 0.08,
        "top_feeder_industries": [],
    }
    return data


def _collect_seasonal_trends(industry: str) -> Dict[str, Any]:
    """Monthly hiring seasonality patterns."""
    _lazy_load()
    data: Dict[str, Any] = {"source": "knowledge_base", "monthly": {}}

    month_names = ["", "Jan", "Feb", "Mar", "Apr", "May", "Jun",
                   "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]

    if _trend_engine:
        for m in range(1, 13):
            adj = _safe_get(_trend_engine.get_seasonal_adjustment, "mixed", m, default={})
            if adj and adj.get("multiplier"):
                data["monthly"][m] = {"multiplier": adj["multiplier"], "label": month_names[m]}
                data["source"] = "trend_engine"

    if not data["monthly"]:
        for m in range(1, 13):
            data["monthly"][m] = {"multiplier": _SEASONAL_MULT[m], "label": month_names[m]}

    if _research:
        advice = _safe_get(_research.get_seasonal_hiring_advice, industry, default={})
        if advice:
            data["hiring_advice"] = advice
    if "hiring_advice" not in data:
        data["hiring_advice"] = {
            "peak_months": ["Jan", "Sep"],
            "note": "Standard hiring follows Q1 budget releases and fall planning cycles.",
        }

    sorted_months = sorted(data["monthly"].items(), key=lambda x: x[1]["multiplier"])
    data["trough_months"] = [m for m, _ in sorted_months[:3]]
    data["peak_months"] = [m for m, _ in sorted_months[-3:]]
    return data


# ═══════════════════════════════════════════════════════════════════════════════
# NARRATIVE & RECOMMENDATIONS
# ═══════════════════════════════════════════════════════════════════════════════

def _build_executive_summary(report: Dict[str, Any]) -> str:
    paras: List[str] = []
    ind = _industry_label(report.get("industry", ""))
    role_cat = report.get("role_category", "the target role category")
    locs = ", ".join(report.get("locations", [])[:3]) or "multiple regions"

    paras.append(
        f"This market intelligence report analyzes the recruitment landscape for "
        f"{role_cat.replace('_', ' ')} positions in the {ind} sector across {locs}. "
        f"The analysis covers labor market dynamics, compensation benchmarks, "
        f"advertising channel performance, competitor activity, talent supply, "
        f"and seasonal hiring patterns."
    )

    overview = report.get("market_overview", {})
    diff = overview.get("difficulty_score", 50)
    diff_text = "highly competitive" if diff >= 70 else ("moderately competitive" if diff >= 50 else "relatively accessible")
    paras.append(
        f"The labor market is {diff_text} (difficulty score: {diff}/100). "
        f"Industry growth rate stands at {_fmt_pct(overview.get('growth_rate', 0.05))} "
        f"with average time-to-fill of {overview.get('avg_time_to_fill', 40)} days."
    )

    sal = report.get("salary_benchmarks", {}).get("aggregate", {})
    if sal.get("median"):
        paras.append(
            f"Compensation benchmarks indicate an aggregate median salary of "
            f"{_fmt_currency(sal['median'])} across analyzed locations."
        )

    ch = report.get("channel_performance", {})
    ranked = ch.get("ranked_by_cpc", [])
    if ranked:
        cheapest = ranked[0]
        paras.append(
            f"Among advertising channels, {cheapest[0].replace('_', ' ').title()} "
            f"offers the lowest CPC at {_fmt_currency(cheapest[1]['cpc'], 2)}."
        )

    return "\n\n".join(paras)


def _build_recommendations(report: Dict[str, Any]) -> List[Dict[str, Any]]:
    recs: List[Dict[str, Any]] = []
    ch = report.get("channel_performance", {})
    overview = report.get("market_overview", {})
    sal = report.get("salary_benchmarks", {})
    seasonal = report.get("seasonal_trends", {})
    collar = report.get("collar_type", "white_collar")
    diff = overview.get("difficulty_score", 50)

    ranked = ch.get("ranked_by_cpc", [])
    if ranked:
        c = ranked[0]
        recs.append({
            "title": f"Prioritize {c[0].replace('_', ' ').title()} for cost efficiency",
            "description": f"At {_fmt_currency(c[1]['cpc'], 2)} CPC, this channel offers the best cost-per-click. Allocate 30-40% of budget here.",
            "impact": "high", "impact_score": 9, "category": "channel",
        })

    trough = seasonal.get("trough_months", [])
    if trough:
        mn = {1: "Jan", 2: "Feb", 3: "Mar", 4: "Apr", 5: "May", 6: "Jun",
              7: "Jul", 8: "Aug", 9: "Sep", 10: "Oct", 11: "Nov", 12: "Dec"}
        labels = [mn.get(m, str(m)) for m in trough]
        recs.append({
            "title": "Launch campaigns during low-competition months",
            "description": f"CPCs are typically lowest in {', '.join(labels)}. Front-loading budget can reduce CPA by 10-15%.",
            "impact": "medium", "impact_score": 7, "category": "timing",
        })

    if diff >= 65:
        recs.append({
            "title": "Highlight compensation to attract scarce talent",
            "description": f"Difficulty score {diff}/100 -- prominently featuring salary ranges increases apply rates by 20-30%.",
            "impact": "high", "impact_score": 8, "category": "compensation",
        })

    if collar == "blue_collar":
        recs.append({
            "title": "Use mobile-first and SMS-based application flows",
            "description": "Blue-collar candidates are 3x more likely to apply via mobile. Use Indeed and programmatic boards as primary channels.",
            "impact": "high", "impact_score": 8, "category": "strategy",
        })
    elif collar == "grey_collar":
        recs.append({
            "title": "Target niche healthcare/clinical boards",
            "description": "Grey-collar clinical roles convert best on specialty boards (Health eCareers, Nurse.com) alongside Indeed.",
            "impact": "high", "impact_score": 8, "category": "strategy",
        })
    else:
        recs.append({
            "title": "Leverage LinkedIn and employer branding",
            "description": "White-collar candidates respond to employer brand. Invest in LinkedIn Sponsored Jobs and company page content.",
            "impact": "medium", "impact_score": 6, "category": "strategy",
        })

    if len(ranked) >= 3:
        recs.append({
            "title": "Diversify across 3-4 channels for resilience",
            "description": "A diversified mix across job boards, search, social, and programmatic reduces risk from CPC spikes.",
            "impact": "medium", "impact_score": 6, "category": "budget",
        })

    if diff >= 75:
        recs.append({
            "title": "Consider passive candidate targeting",
            "description": f"Difficulty {diff}/100 indicates a very tight market. Supplement with LinkedIn InMail and retargeting.",
            "impact": "high", "impact_score": 9, "category": "strategy",
        })

    comps = report.get("competitor_analysis", {}).get("competitors", [])
    if comps:
        recs.append({
            "title": "Differentiate from competitor employer brands",
            "description": f"With {len(comps)} active competitors, unique EVP messaging and faster apply flows provide an edge.",
            "impact": "medium", "impact_score": 7, "category": "strategy",
        })

    recs.sort(key=lambda r: r.get("impact_score", 0), reverse=True)
    return recs


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN ENTRY POINT
# ═══════════════════════════════════════════════════════════════════════════════

def generate_report(data: Dict[str, Any]) -> Dict[str, Any]:
    """Generate a comprehensive market intelligence report.

    Parameters
    ----------
    data : dict
        - industry (str): e.g. "technology", "healthcare_medical"
        - role_category (str): e.g. "software_engineering", "nursing"
        - locations (list[str]): geographic regions
        - time_period (str, optional): "quarterly" or "annual"
        - competitors (list[str], optional): competitor company names

    Returns
    -------
    dict with keys: market_overview, salary_benchmarks, channel_performance,
        competitor_analysis, talent_supply, seasonal_trends, recommendations,
        collar_type, report_metadata
    """
    _lazy_load()
    industry = data.get("industry", "technology")
    role_category = data.get("role_category", "general")
    locations = data.get("locations", ["United States"])
    time_period = data.get("time_period", "quarterly")
    competitors = data.get("competitors", [])

    collar_type = _detect_collar(role_category, industry)

    report: Dict[str, Any] = {
        "industry": industry,
        "role_category": role_category,
        "locations": locations,
        "collar_type": collar_type,
    }

    # Parallel data collection
    with ThreadPoolExecutor(max_workers=5) as pool:
        futures = {
            "market_overview": pool.submit(_collect_market_overview, industry, locations),
            "salary_benchmarks": pool.submit(_collect_salary_benchmarks, role_category, industry, locations),
            "channel_performance": pool.submit(_collect_channel_performance, industry, collar_type),
            "talent_supply": pool.submit(_collect_talent_supply, role_category, industry, locations),
            "seasonal_trends": pool.submit(_collect_seasonal_trends, industry),
        }
        if competitors:
            futures["competitor_analysis"] = pool.submit(_collect_competitor_analysis, competitors, industry)

        sources: List[Dict[str, str]] = []
        confidence_scores: Dict[str, str] = {}
        for key, fut in futures.items():
            try:
                result = fut.result(timeout=30)
                report[key] = result
                src = result.get("source", "estimate")
                sources.append({"section": key, "source": src})
                confidence_scores[key] = "high" if "api" in src else ("medium" if src != "estimate" else "low")
            except Exception as exc:
                logger.error("Collection failed for %s: %s", key, exc)
                fallback = {"source": "error", "error": str(exc)}
                if key == "competitor_analysis":
                    fallback["competitors"] = []
                elif key == "channel_performance":
                    fallback["channels"] = {}
                    fallback["ranked_by_cpc"] = []
                    fallback["ranked_by_cpa"] = []
                elif key == "seasonal_trends":
                    fallback["monthly"] = {}
                    fallback["trough_months"] = []
                    fallback["peak_months"] = []
                elif key == "salary_benchmarks":
                    fallback["by_location"] = {}
                    fallback["aggregate"] = {}
                elif key == "market_overview":
                    fallback["difficulty_score"] = 50
                report[key] = fallback
                sources.append({"section": key, "source": "error"})
                confidence_scores[key] = "none"

    if not competitors:
        report["competitor_analysis"] = {"source": "none", "competitors": [], "note": "No competitors specified."}

    report["recommendations"] = _build_recommendations(report)
    report["executive_summary"] = _build_executive_summary(report)

    report["report_metadata"] = {
        "generated_at": datetime.now().isoformat(),
        "time_period": time_period,
        "data_sources": sources,
        "confidence_scores": confidence_scores,
        "industry_label": _industry_label(industry),
        "collar_type": collar_type,
        "locations_analyzed": len(locations),
        "competitors_analyzed": len(competitors),
    }

    return report


# ═══════════════════════════════════════════════════════════════════════════════
# EXCEL EXPORT  (Sapphire Blue: #0F172A / #2563EB / #DBEAFE, Calibri, col B)
# ═══════════════════════════════════════════════════════════════════════════════

_XL_DARK = "0F172A"
_XL_BLUE = "2563EB"
_XL_LIGHT = "DBEAFE"
_XL_WHITE = "FFFFFF"
_XL_FONT = "Calibri"


def generate_intel_excel(report_data: Dict[str, Any]) -> bytes:
    """Generate an Excel workbook from report data. Returns xlsx bytes."""
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    wb = Workbook()
    dark_fill = PatternFill("solid", fgColor=_XL_DARK)
    blue_fill = PatternFill("solid", fgColor=_XL_BLUE)
    light_fill = PatternFill("solid", fgColor=_XL_LIGHT)
    white_fill = PatternFill("solid", fgColor=_XL_WHITE)
    title_font = Font(name=_XL_FONT, size=16, bold=True, color=_XL_WHITE)
    hdr_font = Font(name=_XL_FONT, size=11, bold=True, color=_XL_WHITE)
    sub_font = Font(name=_XL_FONT, size=11, bold=True, color=_XL_DARK)
    body_font = Font(name=_XL_FONT, size=10, color="333333")
    num_font = Font(name=_XL_FONT, size=10, color="1E293B")
    thin_border = Border(bottom=Side(style="thin", color="CBD5E1"))
    c_align = Alignment(horizontal="center", vertical="center")
    l_align = Alignment(horizontal="left", vertical="center", wrap_text=True)

    def _hdr_row(ws, row, c1, c2):
        for c in range(c1, c2 + 1):
            cell = ws.cell(row=row, column=c)
            cell.fill = blue_fill
            cell.font = hdr_font
            cell.alignment = c_align

    def _data_row(ws, row, c1, c2, alt=False):
        for c in range(c1, c2 + 1):
            cell = ws.cell(row=row, column=c)
            cell.fill = light_fill if alt else white_fill
            cell.font = body_font
            cell.border = thin_border
            cell.alignment = l_align

    ind_label = _industry_label(report_data.get("industry", ""))
    role_cat = report_data.get("role_category", "N/A").replace("_", " ").title()
    locs = ", ".join(report_data.get("locations", [])[:5])
    meta = report_data.get("report_metadata", {})

    # -- Sheet 1: Overview --
    ws = wb.active
    ws.title = "Overview"
    ws.sheet_properties.tabColor = _XL_BLUE
    ws.column_dimensions["A"].width = 3
    ws.column_dimensions["B"].width = 28
    ws.column_dimensions["C"].width = 50
    ws.column_dimensions["D"].width = 22

    ws.merge_cells("B2:D2")
    t = ws.cell(row=2, column=2, value="Market Intelligence Report")
    t.font = title_font
    t.fill = dark_fill
    t.alignment = c_align
    for c in range(2, 5):
        ws.cell(row=2, column=c).fill = dark_fill

    info = [
        ("Role Category", role_cat),
        ("Industry", ind_label),
        ("Locations", locs),
        ("Collar Type", report_data.get("collar_type", "N/A").replace("_", " ").title()),
        ("Generated", meta.get("generated_at", datetime.now().isoformat())[:16]),
        ("Difficulty Score", f"{report_data.get('market_overview', {}).get('difficulty_score', 'N/A')}/100"),
    ]
    for i, (lbl, val) in enumerate(info):
        ws.cell(row=4 + i, column=2, value=lbl).font = sub_font
        ws.cell(row=4 + i, column=3, value=val).font = body_font

    summary = report_data.get("executive_summary", "")
    if summary:
        ws.merge_cells("B12:D12")
        ws.cell(row=12, column=2, value="Executive Summary").font = Font(name=_XL_FONT, size=13, bold=True, color=_XL_DARK)
        ws.merge_cells("B13:D22")
        sc = ws.cell(row=13, column=2, value=summary)
        sc.font = body_font
        sc.alignment = Alignment(wrap_text=True, vertical="top")

    # -- Sheet 2: Channel Performance --
    ws2 = wb.create_sheet("Channel Performance")
    ws2.sheet_properties.tabColor = _XL_BLUE
    ws2.column_dimensions["A"].width = 3
    for col_letter, w in [("B", 20), ("C", 12), ("D", 12), ("E", 12), ("F", 14), ("G", 14), ("H", 12)]:
        ws2.column_dimensions[col_letter].width = w

    ws2.merge_cells("B2:H2")
    ws2.cell(row=2, column=2, value="Channel Performance Benchmarks").font = Font(name=_XL_FONT, size=14, bold=True, color=_XL_WHITE)
    for c in range(2, 9):
        ws2.cell(row=2, column=c).fill = dark_fill

    headers = ["Channel", "CPC", "CPA", "CTR", "Conv Rate", "Confidence", "Trend"]
    for i, h in enumerate(headers):
        ws2.cell(row=4, column=2 + i, value=h)
    _hdr_row(ws2, 4, 2, 8)

    channels = report_data.get("channel_performance", {}).get("channels", {})
    for idx, (name, m) in enumerate(channels.items()):
        r = 5 + idx
        ws2.cell(row=r, column=2, value=name.replace("_", " ").title())
        ws2.cell(row=r, column=3, value=f"${m.get('cpc', 0):.2f}")
        cpa = m.get("cpa")
        ws2.cell(row=r, column=4, value=f"${cpa:.2f}" if cpa else "N/A")
        ctr = m.get("ctr")
        ws2.cell(row=r, column=5, value=f"{ctr*100:.1f}%" if ctr else "N/A")
        cr = m.get("conv_rate")
        ws2.cell(row=r, column=6, value=f"{cr*100:.1f}%" if cr else "N/A")
        ws2.cell(row=r, column=7, value=m.get("confidence", "N/A"))
        ws2.cell(row=r, column=8, value=m.get("trend", "N/A"))
        _data_row(ws2, r, 2, 8, alt=(idx % 2 == 1))

    # -- Sheet 3: Salary Benchmarks --
    ws3 = wb.create_sheet("Salary Benchmarks")
    ws3.sheet_properties.tabColor = _XL_BLUE
    ws3.column_dimensions["A"].width = 3
    for col_letter, w in [("B", 22), ("C", 14), ("D", 14), ("E", 14), ("F", 14), ("G", 14), ("H", 10)]:
        ws3.column_dimensions[col_letter].width = w

    ws3.merge_cells("B2:H2")
    ws3.cell(row=2, column=2, value="Salary Benchmarks by Location").font = Font(name=_XL_FONT, size=14, bold=True, color=_XL_WHITE)
    for c in range(2, 9):
        ws3.cell(row=2, column=c).fill = dark_fill

    headers = ["Location", "P10", "P25", "Median", "P75", "P90", "COLI"]
    for i, h in enumerate(headers):
        ws3.cell(row=4, column=2 + i, value=h)
    _hdr_row(ws3, 4, 2, 8)

    by_loc = report_data.get("salary_benchmarks", {}).get("by_location", {})
    for idx, (loc, sal) in enumerate(by_loc.items()):
        r = 5 + idx
        ws3.cell(row=r, column=2, value=loc)
        ws3.cell(row=r, column=3, value=_fmt_currency(sal.get("p10", 0)))
        ws3.cell(row=r, column=4, value=_fmt_currency(sal.get("p25", 0)))
        ws3.cell(row=r, column=5, value=_fmt_currency(sal.get("median", 0)))
        ws3.cell(row=r, column=6, value=_fmt_currency(sal.get("p75", 0)))
        ws3.cell(row=r, column=7, value=_fmt_currency(sal.get("p90", 0)))
        ws3.cell(row=r, column=8, value=sal.get("coli", "N/A"))
        _data_row(ws3, r, 2, 8, alt=(idx % 2 == 1))

    # -- Sheet 4: Seasonal Patterns --
    ws4 = wb.create_sheet("Seasonal Patterns")
    ws4.sheet_properties.tabColor = _XL_BLUE
    ws4.column_dimensions["A"].width = 3
    ws4.column_dimensions["B"].width = 14
    ws4.column_dimensions["C"].width = 18
    ws4.column_dimensions["D"].width = 30

    ws4.merge_cells("B2:D2")
    ws4.cell(row=2, column=2, value="Seasonal Hiring Patterns").font = Font(name=_XL_FONT, size=14, bold=True, color=_XL_WHITE)
    for c in range(2, 5):
        ws4.cell(row=2, column=c).fill = dark_fill

    for i, h in enumerate(["Month", "Multiplier", "Interpretation"]):
        ws4.cell(row=4, column=2 + i, value=h)
    _hdr_row(ws4, 4, 2, 4)

    month_names = ["", "January", "February", "March", "April", "May", "June",
                   "July", "August", "September", "October", "November", "December"]
    mults = report_data.get("seasonal_trends", {}).get("monthly", {})
    for idx in range(12):
        mn = idx + 1
        r = 5 + idx
        info = mults.get(mn, mults.get(str(mn), {}))
        mult = info.get("multiplier", 1.0) if isinstance(info, dict) else 1.0
        interp = "Above average (higher competition)" if mult > 1.05 else ("Below average (opportunity)" if mult < 0.95 else "Average")
        ws4.cell(row=r, column=2, value=month_names[mn])
        ws4.cell(row=r, column=3, value=f"{mult:.2f}x")
        ws4.cell(row=r, column=4, value=interp)
        _data_row(ws4, r, 2, 4, alt=(idx % 2 == 1))

    # -- Sheet 5: Recommendations --
    ws5 = wb.create_sheet("Recommendations")
    ws5.sheet_properties.tabColor = _XL_BLUE
    ws5.column_dimensions["A"].width = 3
    for col_letter, w in [("B", 6), ("C", 36), ("D", 60), ("E", 10), ("F", 14)]:
        ws5.column_dimensions[col_letter].width = w

    ws5.merge_cells("B2:F2")
    ws5.cell(row=2, column=2, value="Strategic Recommendations").font = Font(name=_XL_FONT, size=14, bold=True, color=_XL_WHITE)
    for c in range(2, 7):
        ws5.cell(row=2, column=c).fill = dark_fill

    for i, h in enumerate(["#", "Recommendation", "Details", "Impact", "Category"]):
        ws5.cell(row=4, column=2 + i, value=h)
    _hdr_row(ws5, 4, 2, 6)

    for idx, rec in enumerate(report_data.get("recommendations", [])):
        r = 5 + idx
        ws5.cell(row=r, column=2, value=idx + 1)
        ws5.cell(row=r, column=3, value=rec.get("title", ""))
        ws5.cell(row=r, column=4, value=rec.get("description", ""))
        ws5.cell(row=r, column=5, value=rec.get("impact", "").upper())
        ws5.cell(row=r, column=6, value=rec.get("category", ""))
        _data_row(ws5, r, 2, 6, alt=(idx % 2 == 1))

    # -- Sheet 6: Competitor Analysis (if present) --
    comps = report_data.get("competitor_analysis", {}).get("competitors", [])
    if comps:
        ws6 = wb.create_sheet("Competitors")
        ws6.sheet_properties.tabColor = _XL_BLUE
        ws6.column_dimensions["A"].width = 3
        for col_letter, w in [("B", 24), ("C", 18), ("D", 30), ("E", 16)]:
            ws6.column_dimensions[col_letter].width = w

        ws6.merge_cells("B2:E2")
        ws6.cell(row=2, column=2, value="Competitor Analysis").font = Font(name=_XL_FONT, size=14, bold=True, color=_XL_WHITE)
        for c in range(2, 6):
            ws6.cell(row=2, column=c).fill = dark_fill

        for i, h in enumerate(["Company", "Est. Openings", "Primary Channels", "Brand Strength"]):
            ws6.cell(row=4, column=2 + i, value=h)
        _hdr_row(ws6, 4, 2, 5)

        for idx, comp in enumerate(comps):
            r = 5 + idx
            ws6.cell(row=r, column=2, value=comp.get("name", ""))
            ws6.cell(row=r, column=3, value=comp.get("estimated_openings", "N/A"))
            ws6.cell(row=r, column=4, value=", ".join(comp.get("primary_channels", [])))
            ws6.cell(row=r, column=5, value=comp.get("employer_brand_strength", "N/A"))
            _data_row(ws6, r, 2, 5, alt=(idx % 2 == 1))

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.read()


# ═══════════════════════════════════════════════════════════════════════════════
# PPT EXPORT  (Port Gore #202058, Blue Violet #5A54BD, Downy teal #6BB3CD)
# ═══════════════════════════════════════════════════════════════════════════════

def generate_intel_ppt(report_data: Dict[str, Any]) -> bytes:
    """Generate a PowerPoint deck from report data. Returns pptx bytes."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    port_gore = RGBColor(0x20, 0x20, 0x58)
    blue_violet = RGBColor(0x5A, 0x54, 0xBD)
    downy = RGBColor(0x6B, 0xB3, 0xCD)
    white = RGBColor(0xFF, 0xFF, 0xFF)
    light_bg = RGBColor(0xF0, 0xF0, 0xF8)
    dark_text = RGBColor(0x1E, 0x1E, 0x2E)
    muted = RGBColor(0x64, 0x74, 0x8B)

    ind_label = _industry_label(report_data.get("industry", ""))
    role_cat = report_data.get("role_category", "N/A").replace("_", " ").title()
    locs = ", ".join(report_data.get("locations", [])[:3])

    def _bg(slide, color):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = color

    def _tb(slide, left, top, w, h, text, sz=12, bold=False, color=dark_text, align=PP_ALIGN.LEFT):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(text)
        p.font.size = Pt(sz)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = align
        return box

    def _rect(slide, left, top, w, h, fill_color):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        return shape

    # -- Slide 1: Title --
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s1, port_gore)
    _rect(s1, 0, 2.5, 13.333, 0.08, downy)
    _tb(s1, 1, 1.5, 11, 1.2, "MARKET INTELLIGENCE REPORT", sz=32, bold=True, color=white, align=PP_ALIGN.CENTER)
    _tb(s1, 1, 3.0, 11, 0.8, f"{role_cat}  |  {ind_label}  |  {locs}", sz=18, color=downy, align=PP_ALIGN.CENTER)
    _tb(s1, 1, 5.5, 11, 0.5, f"Generated {datetime.now().strftime('%B %d, %Y')}  |  Powered by Joveo",
        sz=12, color=white, align=PP_ALIGN.CENTER)

    # -- Slide 2: Executive Summary --
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s2, white)
    _rect(s2, 0, 0, 13.333, 1.1, port_gore)
    _tb(s2, 0.8, 0.2, 11, 0.7, "Executive Summary", sz=24, bold=True, color=white)

    summary = report_data.get("executive_summary", "No summary available.")
    if len(summary) > 1200:
        summary = summary[:1197] + "..."
    _tb(s2, 0.8, 1.5, 11.5, 5.0, summary, sz=13, color=dark_text)

    diff_score = report_data.get("market_overview", {}).get("difficulty_score", 50)
    _rect(s2, 10.5, 1.5, 2.2, 1.2, blue_violet)
    _tb(s2, 10.6, 1.55, 2.0, 0.4, "HIRING DIFFICULTY", sz=9, bold=True, color=white, align=PP_ALIGN.CENTER)
    _tb(s2, 10.6, 1.95, 2.0, 0.7, f"{diff_score}/100", sz=28, bold=True, color=white, align=PP_ALIGN.CENTER)

    # -- Slide 3: Channel Performance --
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s3, white)
    _rect(s3, 0, 0, 13.333, 1.1, port_gore)
    _tb(s3, 0.8, 0.2, 11, 0.7, "Channel Performance Benchmarks", sz=24, bold=True, color=white)

    channels = report_data.get("channel_performance", {}).get("channels", {})
    card_w = 1.85
    gap = 0.2
    for idx, (ch, m) in enumerate(channels.items()):
        if idx >= 6:
            break
        x = 0.8 + idx * (card_w + gap)
        _rect(s3, x, 1.5, card_w, 4.5, light_bg)
        _tb(s3, x + 0.1, 1.6, card_w - 0.2, 0.5, ch.replace("_", " ").title(), sz=11, bold=True, color=port_gore, align=PP_ALIGN.CENTER)
        _rect(s3, x + 0.15, 2.3, card_w - 0.3, 0.9, blue_violet)
        _tb(s3, x + 0.2, 2.35, card_w - 0.4, 0.3, "CPC", sz=9, color=white, align=PP_ALIGN.CENTER)
        _tb(s3, x + 0.2, 2.65, card_w - 0.4, 0.5, f"${m.get('cpc', 0):.2f}", sz=20, bold=True, color=white, align=PP_ALIGN.CENTER)
        cpa = m.get("cpa")
        _tb(s3, x + 0.1, 3.5, card_w - 0.2, 0.3, "CPA", sz=9, bold=True, color=blue_violet, align=PP_ALIGN.CENTER)
        _tb(s3, x + 0.1, 3.8, card_w - 0.2, 0.4, f"${cpa:.2f}" if cpa else "N/A", sz=16, bold=True, color=dark_text, align=PP_ALIGN.CENTER)
        _tb(s3, x + 0.1, 4.5, card_w - 0.2, 0.4, f"Trend: {m.get('trend', 'stable').title()}", sz=10, color=port_gore, align=PP_ALIGN.CENTER)

    # -- Slide 4: Salary Benchmarks --
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s4, white)
    _rect(s4, 0, 0, 13.333, 1.1, port_gore)
    _tb(s4, 0.8, 0.2, 11, 0.7, "Salary Benchmarks by Location", sz=24, bold=True, color=white)

    by_loc = report_data.get("salary_benchmarks", {}).get("by_location", {})
    y_off = 1.5
    for idx, (loc, sal) in enumerate(list(by_loc.items())[:6]):
        y = y_off + idx * 0.9
        _tb(s4, 1.0, y, 2.5, 0.4, loc, sz=11, bold=True, color=dark_text)
        med = sal.get("median", 0)
        max_sal = max((s.get("median", 0) for s in by_loc.values()), default=1) or 1
        bar_w = max(0.5, (med / max_sal) * 7.0)
        _rect(s4, 3.5, y + 0.05, bar_w, 0.35, downy)
        _tb(s4, 3.6 + bar_w, y, 2.0, 0.4, _fmt_currency(med), sz=11, bold=True, color=dark_text)

    agg = report_data.get("salary_benchmarks", {}).get("aggregate", {})
    if agg:
        _rect(s4, 1.0, 7.0, 11.3, 0.3, light_bg)
        _tb(s4, 1.2, 7.0, 10.0, 0.3,
            f"Aggregate median: {_fmt_currency(agg.get('median', 0))}  |  Range: {_fmt_currency(agg.get('min_p10', 0))} - {_fmt_currency(agg.get('max_p90', 0))}",
            sz=10, bold=True, color=port_gore)

    # -- Slide 5: Seasonal Patterns --
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s5, white)
    _rect(s5, 0, 0, 13.333, 1.1, port_gore)
    _tb(s5, 0.8, 0.2, 11, 0.7, "Seasonal Hiring Patterns", sz=24, bold=True, color=white)

    mults = report_data.get("seasonal_trends", {}).get("monthly", {})
    month_abbrs = ["Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]
    bar_base = 5.5
    max_mult = max((info.get("multiplier", 1.0) if isinstance(info, dict) else 1.0 for info in mults.values()), default=1.2)

    for idx in range(12):
        mn = idx + 1
        info = mults.get(mn, mults.get(str(mn), {}))
        mult = info.get("multiplier", 1.0) if isinstance(info, dict) else 1.0
        x = 0.8 + idx * 1.02
        bar_h = max(0.2, (mult / max_mult) * 3.5)
        y = bar_base - bar_h
        color = downy if mult < 0.97 else (blue_violet if mult > 1.03 else RGBColor(0x94, 0xA3, 0xB8))
        _rect(s5, x, y, 0.85, bar_h, color)
        _tb(s5, x, bar_base + 0.1, 0.85, 0.3, month_abbrs[idx], sz=9, bold=True, color=dark_text, align=PP_ALIGN.CENTER)
        _tb(s5, x, y - 0.3, 0.85, 0.3, f"{mult:.2f}x", sz=8, color=dark_text, align=PP_ALIGN.CENTER)

    advice = report_data.get("seasonal_trends", {}).get("hiring_advice", {})
    note = advice.get("note", "")
    if note:
        _tb(s5, 0.8, 6.3, 11.5, 0.8, note, sz=10, color=port_gore)

    # -- Slide 6: Recommendations --
    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s6, white)
    _rect(s6, 0, 0, 13.333, 1.1, port_gore)
    _tb(s6, 0.8, 0.2, 11, 0.7, "Key Recommendations", sz=24, bold=True, color=white)

    recs = report_data.get("recommendations", [])
    for idx, rec in enumerate(recs[:6]):
        y = 1.4 + idx * 0.95
        impact = rec.get("impact", "medium")
        badge_c = blue_violet if impact == "high" else (downy if impact == "medium" else RGBColor(0x94, 0xA3, 0xB8))
        _rect(s6, 0.8, y, 0.15, 0.7, badge_c)
        _tb(s6, 1.1, y, 4.0, 0.35, rec.get("title", ""), sz=11, bold=True, color=dark_text)
        _tb(s6, 1.1, y + 0.35, 10.5, 0.35, rec.get("description", ""), sz=9, color=muted)
        _tb(s6, 11.8, y + 0.05, 1.0, 0.3, impact.upper(), sz=9, bold=True, color=badge_c, align=PP_ALIGN.RIGHT)

    # -- Slide 7: Closing --
    s7 = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s7, port_gore)
    _rect(s7, 0, 3.5, 13.333, 0.06, downy)
    _tb(s7, 1, 2.5, 11, 0.8, "Thank You", sz=36, bold=True, color=white, align=PP_ALIGN.CENTER)
    _tb(s7, 1, 4.0, 11, 0.6, "Powered by Joveo  |  Intelligent Recruitment Advertising",
        sz=14, color=downy, align=PP_ALIGN.CENTER)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ═══════════════════════════════════════════════════════════════════════════════
# UNIFIED HTTP HANDLER
# ═══════════════════════════════════════════════════════════════════════════════

def handle_market_intel_request(
    path: str, method: str, body: Optional[Dict[str, Any]] = None,
):
    """Unified handler for /api/market-intel/* routes.

    Routes:
        POST /api/market-intel/generate  -> generate_report(body) -> dict
        POST /api/market-intel/excel     -> generate_intel_excel(body) -> bytes
        POST /api/market-intel/ppt       -> generate_intel_ppt(body) -> bytes

    Returns:
        - bytes for excel/ppt routes (app.py sends as binary download)
        - dict for generate route (app.py sends as JSON)
    """
    body = body or {}

    try:
        # POST /api/market-intel/generate
        if path == "/api/market-intel/generate" and method == "POST":
            report = generate_report(body)
            # Strip bytes from JSON response
            json_safe = {k: v for k, v in report.items() if not isinstance(v, bytes)}
            return {"ok": True, "report": json_safe}

        # POST /api/market-intel/excel
        if path == "/api/market-intel/excel" and method == "POST":
            if "market_overview" in body:
                report_data = body
            else:
                report_data = generate_report(body)
            return generate_intel_excel(report_data)

        # POST /api/market-intel/ppt
        if path == "/api/market-intel/ppt" and method == "POST":
            if "market_overview" in body:
                report_data = body
            else:
                report_data = generate_report(body)
            return generate_intel_ppt(report_data)

        return {"ok": False, "error": f"Unknown route: {method} {path}"}

    except Exception as exc:
        logger.exception("market_intel_request error on %s %s", method, path)
        return {"ok": False, "error": str(exc)}


# ═══════════════════════════════════════════════════════════════════════════════
# BACKWARD COMPATIBILITY -- old API surface
# ═══════════════════════════════════════════════════════════════════════════════

def generate_market_intel_report(
    role: str, industry: str, location: str,
    options: Optional[Dict[str, Any]] = None,
) -> Dict[str, Any]:
    """Legacy entry point -- wraps generate_report for backward compatibility."""
    options = options or {}
    report = generate_report({
        "industry": industry,
        "role_category": role,
        "locations": [location],
    })

    # Map new keys to old keys for compatibility
    report["role"] = role
    report["location"] = location
    report["metadata"] = report.get("report_metadata", {})

    # Compensation mapping
    sal = report.get("salary_benchmarks", {})
    agg = sal.get("aggregate", {})
    by_loc = sal.get("by_location", {})
    first_loc = next(iter(by_loc.values()), {})
    report["compensation"] = {
        "salary_ranges": {
            "p10": first_loc.get("p10", agg.get("min_p10", 35000)),
            "p25": first_loc.get("p25", 45000),
            "median": agg.get("median", 58000),
            "p75": first_loc.get("p75", 78000),
            "p90": first_loc.get("p90", agg.get("max_p90", 105000)),
            "source": sal.get("source", "knowledge_base"),
        },
        "source": sal.get("source", "knowledge_base"),
    }

    # Labor market mapping
    overview = report.get("market_overview", {})
    report["labor_market"] = {
        "difficulty_score": overview.get("difficulty_score", 50),
        "supply_demand_ratio": 0.7 if overview.get("talent_supply") == "tight" else 1.2,
        "avg_time_to_fill_days": overview.get("avg_time_to_fill", 40),
        "source": overview.get("source", "knowledge_base"),
    }

    # Seasonal mapping
    seasonal = report.get("seasonal_trends", {})
    report["seasonal_patterns"] = {
        "monthly_multipliers": seasonal.get("monthly", {}),
        "hiring_advice": seasonal.get("hiring_advice", {}),
        "peak_months": seasonal.get("peak_months", []),
        "trough_months": seasonal.get("trough_months", []),
    }

    # CPC trends stub
    report["cpc_trends"] = {
        "platform_trends": {},
        "aggregate_trend": {"avg_yoy_change": -3.0, "direction": "declining"},
        "source": "estimate",
    }

    # Generate exports if requested
    if options.get("include_excel", True):
        try:
            report["excel_bytes"] = generate_intel_excel(report)
        except Exception as exc:
            logger.error("Excel generation failed: %s", exc)
            report["excel_bytes"] = None

    if options.get("include_ppt", True):
        try:
            report["ppt_bytes"] = generate_intel_ppt(report)
        except Exception as exc:
            logger.error("PPT generation failed: %s", exc)
            report["ppt_bytes"] = None

    return report


# ═══════════════════════════════════════════════════════════════════════════════
# CLI demo
# ═══════════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    import sys as _sys

    logging.basicConfig(level=logging.INFO)

    industry = _sys.argv[1] if len(_sys.argv) > 1 else "technology"
    role_cat = _sys.argv[2] if len(_sys.argv) > 2 else "software_engineering"
    locs = _sys.argv[3].split(",") if len(_sys.argv) > 3 else ["San Francisco", "New York"]

    print(f"Generating market intel report: {role_cat} / {industry} / {locs}")
    result = generate_report({
        "industry": industry,
        "role_category": role_cat,
        "locations": locs,
        "competitors": ["Google", "Microsoft"],
    })

    output = {k: v for k, v in result.items() if not isinstance(v, bytes)}
    print(json.dumps(output, indent=2, default=str))
