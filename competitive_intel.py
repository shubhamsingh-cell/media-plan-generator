#!/usr/bin/env python3
"""
competitive_intel.py -- Competitive Intelligence Dashboard Backend

Provides comprehensive competitive analysis for recruitment advertising:
  - Company profiling via Wikipedia/Clearbit APIs
  - Multi-competitor comparison (size, industry, public/private)
  - Hiring activity & difficulty benchmarks
  - Ad platform CPC/CPA benchmarks across industries
  - Google Trends search interest comparison
  - Excel & PowerPoint report generation

Thread-safe, graceful degradation when APIs are unavailable.
All external API calls use ThreadPoolExecutor for concurrent fetching.

Depends on (lazy-imported):
  - api_enrichment (fetch_company_info, fetch_company_metadata, fetch_search_trends, fetch_sec_data)
  - data_orchestrator (enrich_competitive, enrich_employer_brand)
  - research (get_competitors, get_company_intelligence, get_client_competitor_intelligence)
  - trend_engine (get_all_platform_benchmarks, get_benchmark, PLATFORMS, INDUSTRIES)
  - shared_utils (INDUSTRY_LABEL_MAP)
"""

from __future__ import annotations

import io
import json
import logging
import re
import threading
import time
from concurrent.futures import ThreadPoolExecutor, as_completed
from datetime import datetime
from typing import Any, Dict, List, Optional

logger = logging.getLogger(__name__)


# ═══════════════════════════════════════════════════════════════════════════════
# LLM ROUTER (lazy import for competitive narrative intelligence)
# ═══════════════════════════════════════════════════════════════════════════════
_llm_router_ci = None
_llm_router_ci_checked = False
_llm_router_ci_lock = threading.Lock()


def _lazy_llm_router_ci():
    """Lazy-load LLM router for competitive narrative generation.

    Returns the module or None if unavailable. Thread-safe.
    """
    global _llm_router_ci, _llm_router_ci_checked
    if _llm_router_ci_checked:
        return _llm_router_ci
    with _llm_router_ci_lock:
        if _llm_router_ci_checked:
            return _llm_router_ci
        try:
            import llm_router as _mod

            _llm_router_ci = _mod
        except ImportError:
            logger.warning("llm_router not available; competitive narrative disabled")
            _llm_router_ci = None
        _llm_router_ci_checked = True
    return _llm_router_ci


def _generate_competitive_narrative(competitor_data: Dict[str, Any]) -> str:
    """Generate AI narrative synthesizing the competitive landscape.

    Args:
        competitor_data: Aggregated competitor analysis results.

    Returns:
        Strategic assessment string, or empty string on failure.
    """
    router = _lazy_llm_router_ci()
    if not router:
        return ""
    # S50: Route competitive summaries to TASK_INTELLIGENCE_SUMMARY (Gemini Flash
    # Lite -- free, fast, ideal for short 3-sentence summaries) instead of
    # TASK_RESEARCH (Claude Haiku -- expensive, overkill for summary length).
    # 10s timeout for non-blocking plan generation.
    task_type = getattr(router, "TASK_INTELLIGENCE_SUMMARY", "intelligence_summary")
    try:
        # S27: Increased from 2000 to 5000 chars to preserve competitor data integrity
        data_snapshot_raw = json.dumps(competitor_data, indent=2, default=str)
        # S27: Smart truncation -- find last complete JSON boundary instead of hard cut
        if len(data_snapshot_raw) > 5000:
            # Find the last closing brace/bracket before the 5000 char limit
            _trunc = data_snapshot_raw[:5000]
            _last_brace = max(_trunc.rfind("}"), _trunc.rfind("]"))
            if _last_brace > 100:
                data_snapshot = _trunc[: _last_brace + 1]
            else:
                data_snapshot = _trunc
        else:
            data_snapshot = data_snapshot_raw
        prompt = (
            f"Synthesize this competitive hiring landscape:\n{data_snapshot}\n\n"
            f"Write a 3-sentence strategic assessment: who's hiring most aggressively, "
            f"where the gaps are, and what this client should do differently."
        )
        result = router.call_llm(
            messages=[{"role": "user", "content": prompt}],
            system_prompt=(
                "You are a senior recruitment marketing strategist specializing in "
                "competitive intelligence. Write concise, actionable assessments. "
                "Cite specific data points. No fluff."
            ),
            task_type=task_type,
            max_tokens=300,
            timeout_budget=10.0,
        )
        return result.get("text") or ""
    except Exception as e:
        logger.error("Competitive narrative generation failed: %s", e, exc_info=True)
        return ""


# ═══════════════════════════════════════════════════════════════════════════════
# Lazy Imports -- graceful fallback when modules are unavailable
# ═══════════════════════════════════════════════════════════════════════════════

_api_enrichment = None
_data_orchestrator = None
_research = None
_trend_engine = None

_HAS_API = False
_HAS_ORCHESTRATOR = False
_HAS_RESEARCH = False
_HAS_TRENDS = False
_HAS_BENCHMARK_REGISTRY = False

try:
    from benchmark_registry import get_channel_benchmark

    _HAS_BENCHMARK_REGISTRY = True
except ImportError:
    _HAS_BENCHMARK_REGISTRY = False

_import_lock = threading.Lock()


def _lazy_api():
    global _api_enrichment, _HAS_API
    if _api_enrichment is not None:
        return _api_enrichment
    with _import_lock:
        if _api_enrichment is not None:
            return _api_enrichment
        try:
            import api_enrichment as _mod

            _api_enrichment = _mod
            _HAS_API = True
        except ImportError:
            logger.warning(
                "api_enrichment not available; company lookups will use fallbacks"
            )
            _HAS_API = False
    return _api_enrichment


def _lazy_orchestrator():
    global _data_orchestrator, _HAS_ORCHESTRATOR
    if _data_orchestrator is not None:
        return _data_orchestrator
    with _import_lock:
        if _data_orchestrator is not None:
            return _data_orchestrator
        try:
            import data_orchestrator as _mod

            _data_orchestrator = _mod
            _HAS_ORCHESTRATOR = True
        except ImportError:
            logger.warning(
                "data_orchestrator not available; competitive enrichment limited"
            )
            _HAS_ORCHESTRATOR = False
    return _data_orchestrator


def _lazy_research():
    global _research, _HAS_RESEARCH
    if _research is not None:
        return _research
    with _import_lock:
        if _research is not None:
            return _research
        try:
            import research as _mod

            _research = _mod
            _HAS_RESEARCH = True
        except ImportError:
            logger.warning("research not available; competitor intelligence limited")
            _HAS_RESEARCH = False
    return _research


def _lazy_trends():
    global _trend_engine, _HAS_TRENDS
    if _trend_engine is not None:
        return _trend_engine
    with _import_lock:
        if _trend_engine is not None:
            return _trend_engine
        try:
            import trend_engine as _mod

            _trend_engine = _mod
            _HAS_TRENDS = True
        except ImportError:
            logger.warning("trend_engine not available; benchmarks will use fallbacks")
            _HAS_TRENDS = False
    return _trend_engine


# ═══════════════════════════════════════════════════════════════════════════════
# Industry label map (import from shared_utils or define fallback)
# ═══════════════════════════════════════════════════════════════════════════════

try:
    from shared_utils import INDUSTRY_LABEL_MAP
except ImportError:
    INDUSTRY_LABEL_MAP = {
        "healthcare_medical": "Healthcare & Medical",
        "blue_collar_trades": "Blue Collar / Skilled Trades",
        "maritime_marine": "Maritime & Marine",
        "military_recruitment": "Military Recruitment",
        "tech_engineering": "Technology & Engineering",
        "general_entry_level": "General / Entry-Level",
        "legal_services": "Legal Services",
        "finance_banking": "Finance & Banking",
        "mental_health": "Mental Health & Behavioral",
        "retail_consumer": "Retail & Consumer",
        "aerospace_defense": "Aerospace & Defense",
        "pharma_biotech": "Pharma & Biotech",
        "energy_utilities": "Energy & Utilities",
        "insurance": "Insurance",
        "telecommunications": "Telecommunications",
        "automotive": "Automotive & Manufacturing",
        "food_beverage": "Food & Beverage",
        "logistics_supply_chain": "Logistics & Supply Chain",
        "hospitality_travel": "Hospitality & Travel",
        "media_entertainment": "Media & Entertainment",
        "construction_real_estate": "Construction & Real Estate",
        "education": "Education",
    }

# Canonical industry keys for dropdown
CANONICAL_INDUSTRIES = list(INDUSTRY_LABEL_MAP.keys())

# Platform display names
PLATFORM_DISPLAY = {
    "google_search": "Google Ads",
    "meta_facebook": "Meta (Facebook)",
    "meta_instagram": "Meta (Instagram)",
    "linkedin": "LinkedIn",
    "indeed": "Indeed",
    "programmatic": "Programmatic",
}

# Max concurrent API workers
_MAX_WORKERS = 10


# ═══════════════════════════════════════════════════════════════════════════════
# WELL-KNOWN COMPANIES FALLBACK DATABASE
# Used when external APIs (Wikipedia, Clearbit, SEC EDGAR) are unavailable or
# return empty results. Covers major employers frequently analyzed.
# ═══════════════════════════════════════════════════════════════════════════════

_WELL_KNOWN_COMPANIES: Dict[str, Dict[str, Any]] = {
    "amazon": {
        "name": "Amazon",
        "description": "Amazon.com, Inc. is an American multinational technology company focusing on e-commerce, cloud computing (AWS), online advertising, digital streaming, and artificial intelligence. It is one of the world's most valuable companies and the largest Internet company by revenue.",
        "industry": "Technology & E-Commerce",
        "employee_count": "1,500,000+",
        "founded": "1994",
        "headquarters": "Seattle, WA",
        "domain": "amazon.com",
        "stock_ticker": "AMZN",
        "is_public": True,
    },
    "google": {
        "name": "Google (Alphabet Inc.)",
        "description": "Google LLC is an American multinational corporation and technology company focusing on search engine technology, online advertising, cloud computing, computer software, quantum computing, e-commerce, and artificial intelligence. It is a subsidiary of Alphabet Inc.",
        "industry": "Technology",
        "employee_count": "180,000+",
        "founded": "1998",
        "headquarters": "Mountain View, CA",
        "domain": "google.com",
        "stock_ticker": "GOOGL",
        "is_public": True,
    },
    "alphabet": {
        "name": "Alphabet Inc.",
        "description": "Alphabet Inc. is an American multinational technology conglomerate and the parent company of Google, YouTube, Waymo, and other subsidiaries. It was created through a restructuring of Google in 2015.",
        "industry": "Technology",
        "employee_count": "180,000+",
        "founded": "2015",
        "headquarters": "Mountain View, CA",
        "domain": "abc.xyz",
        "stock_ticker": "GOOGL",
        "is_public": True,
    },
    "microsoft": {
        "name": "Microsoft Corporation",
        "description": "Microsoft Corporation is an American multinational technology corporation producing computer software, consumer electronics, personal computers, and related services. Its best-known products include Windows, Office, Azure, and Xbox.",
        "industry": "Technology",
        "employee_count": "220,000+",
        "founded": "1975",
        "headquarters": "Redmond, WA",
        "domain": "microsoft.com",
        "stock_ticker": "MSFT",
        "is_public": True,
    },
    "apple": {
        "name": "Apple Inc.",
        "description": "Apple Inc. is an American multinational corporation and technology company headquartered in Cupertino, California. It designs, develops, and sells consumer electronics, computer software, and online services including iPhone, Mac, iPad, and Apple Watch.",
        "industry": "Technology & Consumer Electronics",
        "employee_count": "160,000+",
        "founded": "1976",
        "headquarters": "Cupertino, CA",
        "domain": "apple.com",
        "stock_ticker": "AAPL",
        "is_public": True,
    },
    "meta": {
        "name": "Meta Platforms, Inc.",
        "description": "Meta Platforms, Inc. (formerly Facebook, Inc.) is an American multinational technology conglomerate that owns and operates Facebook, Instagram, WhatsApp, and Threads, and is developing virtual and augmented reality technologies.",
        "industry": "Technology & Social Media",
        "employee_count": "67,000+",
        "founded": "2004",
        "headquarters": "Menlo Park, CA",
        "domain": "meta.com",
        "stock_ticker": "META",
        "is_public": True,
    },
    "facebook": {
        "name": "Meta Platforms, Inc.",
        "description": "Meta Platforms, Inc. (formerly Facebook, Inc.) is an American multinational technology conglomerate that owns and operates Facebook, Instagram, WhatsApp, and Threads.",
        "industry": "Technology & Social Media",
        "employee_count": "67,000+",
        "founded": "2004",
        "headquarters": "Menlo Park, CA",
        "domain": "meta.com",
        "stock_ticker": "META",
        "is_public": True,
    },
    "netflix": {
        "name": "Netflix, Inc.",
        "description": "Netflix, Inc. is an American subscription video on-demand over-the-top streaming service and production company. It offers a library of films and television series through distribution deals and its own productions.",
        "industry": "Media & Entertainment",
        "employee_count": "13,000+",
        "founded": "1997",
        "headquarters": "Los Gatos, CA",
        "domain": "netflix.com",
        "stock_ticker": "NFLX",
        "is_public": True,
    },
    "tesla": {
        "name": "Tesla, Inc.",
        "description": "Tesla, Inc. is an American multinational automotive and clean energy company that designs, manufactures, and sells electric vehicles, battery energy storage, solar panels, and related products and services.",
        "industry": "Automotive & Clean Energy",
        "employee_count": "140,000+",
        "founded": "2003",
        "headquarters": "Austin, TX",
        "domain": "tesla.com",
        "stock_ticker": "TSLA",
        "is_public": True,
    },
    "walmart": {
        "name": "Walmart Inc.",
        "description": "Walmart Inc. is an American multinational retail corporation that operates a chain of hypermarkets, discount department stores, and grocery stores. It is the world's largest company by revenue and the largest private employer.",
        "industry": "Retail",
        "employee_count": "2,100,000+",
        "founded": "1962",
        "headquarters": "Bentonville, AR",
        "domain": "walmart.com",
        "stock_ticker": "WMT",
        "is_public": True,
    },
    "jpmorgan": {
        "name": "JPMorgan Chase & Co.",
        "description": "JPMorgan Chase & Co. is an American multinational financial services firm and the largest bank in the United States by assets. It provides investment banking, financial services, and asset management.",
        "industry": "Finance & Banking",
        "employee_count": "300,000+",
        "founded": "2000",
        "headquarters": "New York, NY",
        "domain": "jpmorganchase.com",
        "stock_ticker": "JPM",
        "is_public": True,
    },
    "jpmorgan chase": {
        "name": "JPMorgan Chase & Co.",
        "description": "JPMorgan Chase & Co. is an American multinational financial services firm and the largest bank in the United States by assets.",
        "industry": "Finance & Banking",
        "employee_count": "300,000+",
        "founded": "2000",
        "headquarters": "New York, NY",
        "domain": "jpmorganchase.com",
        "stock_ticker": "JPM",
        "is_public": True,
    },
    "nike": {
        "name": "Nike, Inc.",
        "description": "Nike, Inc. is an American multinational corporation that designs, develops, manufactures, and sells footwear, apparel, equipment, accessories, and services worldwide. It is the world's largest supplier of athletic shoes and apparel.",
        "industry": "Retail & Consumer Goods",
        "employee_count": "79,000+",
        "founded": "1964",
        "headquarters": "Beaverton, OR",
        "domain": "nike.com",
        "stock_ticker": "NKE",
        "is_public": True,
    },
    "uber": {
        "name": "Uber Technologies, Inc.",
        "description": "Uber Technologies, Inc. is an American multinational transportation company that provides ride-hailing services, courier services, food delivery, and freight transport.",
        "industry": "Technology & Transportation",
        "employee_count": "32,000+",
        "founded": "2009",
        "headquarters": "San Francisco, CA",
        "domain": "uber.com",
        "stock_ticker": "UBER",
        "is_public": True,
    },
    "salesforce": {
        "name": "Salesforce, Inc.",
        "description": "Salesforce, Inc. is an American cloud-based software company that provides customer relationship management (CRM) software and applications focused on sales, customer service, marketing automation, analytics, and application development.",
        "industry": "Technology & SaaS",
        "employee_count": "73,000+",
        "founded": "1999",
        "headquarters": "San Francisco, CA",
        "domain": "salesforce.com",
        "stock_ticker": "CRM",
        "is_public": True,
    },
    "deloitte": {
        "name": "Deloitte Touche Tohmatsu Limited",
        "description": "Deloitte is a multinational professional services network and the largest professional services network in the world by revenue and number of professionals. It provides audit, consulting, financial advisory, risk advisory, tax, and legal services.",
        "industry": "Professional Services",
        "employee_count": "415,000+",
        "founded": "1845",
        "headquarters": "London, UK",
        "domain": "deloitte.com",
        "stock_ticker": "",
        "is_public": False,
    },
    "mckinsey": {
        "name": "McKinsey & Company",
        "description": "McKinsey & Company is an American worldwide management consulting firm that conducts qualitative and quantitative analysis to evaluate management decisions across the public and private sectors.",
        "industry": "Management Consulting",
        "employee_count": "45,000+",
        "founded": "1926",
        "headquarters": "New York, NY",
        "domain": "mckinsey.com",
        "stock_ticker": "",
        "is_public": False,
    },
    "nvidia": {
        "name": "NVIDIA Corporation",
        "description": "NVIDIA Corporation is an American multinational corporation and technology company that designs and supplies graphics processing units (GPUs), application programming interfaces (APIs), and system on a chip units (SoCs) for gaming, professional visualization, data centers, and automotive markets.",
        "industry": "Technology & Semiconductors",
        "employee_count": "30,000+",
        "founded": "1993",
        "headquarters": "Santa Clara, CA",
        "domain": "nvidia.com",
        "stock_ticker": "NVDA",
        "is_public": True,
    },
    "ibm": {
        "name": "International Business Machines Corporation",
        "description": "IBM is an American multinational technology company that produces and sells computer hardware, middleware, and software, and provides hosting and consulting services in areas ranging from mainframe computers to nanotechnology.",
        "industry": "Technology & Consulting",
        "employee_count": "280,000+",
        "founded": "1911",
        "headquarters": "Armonk, NY",
        "domain": "ibm.com",
        "stock_ticker": "IBM",
        "is_public": True,
    },
    "oracle": {
        "name": "Oracle Corporation",
        "description": "Oracle Corporation is an American multinational computer technology company that sells database software and technology, cloud engineered systems, and enterprise software products -- particularly its own brands of database management systems.",
        "industry": "Technology & Enterprise Software",
        "employee_count": "160,000+",
        "founded": "1977",
        "headquarters": "Austin, TX",
        "domain": "oracle.com",
        "stock_ticker": "ORCL",
        "is_public": True,
    },
    "johnson & johnson": {
        "name": "Johnson & Johnson",
        "description": "Johnson & Johnson is an American multinational corporation that develops medical devices, pharmaceuticals, and consumer packaged goods. It is one of the world's most valuable companies and the world's largest healthcare company.",
        "industry": "Healthcare & Pharmaceuticals",
        "employee_count": "130,000+",
        "founded": "1886",
        "headquarters": "New Brunswick, NJ",
        "domain": "jnj.com",
        "stock_ticker": "JNJ",
        "is_public": True,
    },
    "disney": {
        "name": "The Walt Disney Company",
        "description": "The Walt Disney Company is an American multinational mass media and entertainment conglomerate. Its divisions include Disney Entertainment, ESPN, and Disney Experiences, encompassing theme parks, film studios, television networks, and streaming platforms.",
        "industry": "Media & Entertainment",
        "employee_count": "220,000+",
        "founded": "1923",
        "headquarters": "Burbank, CA",
        "domain": "disney.com",
        "stock_ticker": "DIS",
        "is_public": True,
    },
    "boeing": {
        "name": "The Boeing Company",
        "description": "The Boeing Company is an American multinational corporation that designs, manufactures, and sells airplanes, rotorcraft, rockets, satellites, telecommunications equipment, and missiles worldwide.",
        "industry": "Aerospace & Defense",
        "employee_count": "170,000+",
        "founded": "1916",
        "headquarters": "Arlington, VA",
        "domain": "boeing.com",
        "stock_ticker": "BA",
        "is_public": True,
    },
    "goldman sachs": {
        "name": "The Goldman Sachs Group, Inc.",
        "description": "The Goldman Sachs Group, Inc. is an American multinational investment bank and financial services company. It offers services in investment management, securities, asset management, prime brokerage, and securities underwriting.",
        "industry": "Finance & Investment Banking",
        "employee_count": "45,000+",
        "founded": "1869",
        "headquarters": "New York, NY",
        "domain": "goldmansachs.com",
        "stock_ticker": "GS",
        "is_public": True,
    },
    "target": {
        "name": "Target Corporation",
        "description": "Target Corporation is an American retail corporation that operates a chain of discount department stores and hypermarkets. It is the seventh-largest retailer in the United States.",
        "industry": "Retail",
        "employee_count": "400,000+",
        "founded": "1902",
        "headquarters": "Minneapolis, MN",
        "domain": "target.com",
        "stock_ticker": "TGT",
        "is_public": True,
    },
    "indeed": {
        "name": "Indeed",
        "description": "Indeed is an American worldwide employment website for job listings launched in 2004. It is a subsidiary of Recruit Holdings. It aggregates job listings from thousands of websites, including job boards, staffing firms, associations, and company career pages.",
        "industry": "Technology & HR Tech",
        "employee_count": "15,000+",
        "founded": "2004",
        "headquarters": "Austin, TX",
        "domain": "indeed.com",
        "stock_ticker": "",
        "is_public": False,
    },
    "linkedin": {
        "name": "LinkedIn",
        "description": "LinkedIn is a business and employment-focused social media platform owned by Microsoft. It works through websites and mobile apps. It is used for professional networking and career development.",
        "industry": "Technology & Social Media",
        "employee_count": "21,000+",
        "founded": "2002",
        "headquarters": "Sunnyvale, CA",
        "domain": "linkedin.com",
        "stock_ticker": "",
        "is_public": False,
    },
    "joveo": {
        "name": "Joveo",
        "description": "Joveo is a global leader in programmatic job advertising technology. Its AI-powered platform helps enterprises and staffing agencies optimize their recruitment marketing spend across thousands of job sites and channels.",
        "industry": "Technology & HR Tech",
        "employee_count": "200+",
        "founded": "2017",
        "headquarters": "San Francisco, CA",
        "domain": "joveo.com",
        "stock_ticker": "",
        "is_public": False,
    },
}


def _lookup_well_known_company(company_name: str) -> Optional[Dict[str, Any]]:
    """Look up a company in the well-known companies database.

    Performs case-insensitive matching against company names and common
    abbreviations. Returns a copy of the company data or None if not found.

    Args:
        company_name: Company name to look up.

    Returns:
        Dict with company profile data, or None if not in the database.
    """
    if not company_name:
        return None
    key = company_name.strip().lower()
    # Direct match
    if key in _WELL_KNOWN_COMPANIES:
        return dict(_WELL_KNOWN_COMPANIES[key])
    # Try removing common suffixes
    for suffix in (
        " inc",
        " inc.",
        " corp",
        " corp.",
        " llc",
        " ltd",
        " co",
        " company",
        " corporation",
    ):
        stripped = key.rstrip(".").removesuffix(suffix)
        if stripped != key and stripped in _WELL_KNOWN_COMPANIES:
            return dict(_WELL_KNOWN_COMPANIES[stripped])
    # Substring match for names like "The Walt Disney Company" -> "disney"
    for known_key, known_data in _WELL_KNOWN_COMPANIES.items():
        if known_key in key or key in (known_data.get("name") or "").lower():
            return dict(known_data)
    return None


# ═══════════════════════════════════════════════════════════════════════════════
# 1. COMPANY ANALYSIS
# ═══════════════════════════════════════════════════════════════════════════════


def analyze_company(company_name: str) -> Dict[str, Any]:
    """Fetch company profile data using api_enrichment functions.

    Returns dict with: name, description, industry, employee_count,
    founded, headquarters, logo_url, stock_ticker, domain, is_public.
    Gracefully falls back to partial data when APIs are unavailable.
    """
    if not company_name or not company_name.strip():
        return {"name": company_name or "Unknown", "error": "Empty company name"}

    company_name = company_name.strip()
    profile: Dict[str, Any] = {
        "name": company_name,
        "description": "",
        "industry": "",
        "employee_count": None,
        "founded": None,
        "headquarters": "",
        "logo_url": "",
        "stock_ticker": "",
        "domain": "",
        "is_public": False,
        "glassdoor_rating": None,
        "employer_brand_strength": "",
        "sources": [],
    }

    api = _lazy_api()
    orch = _lazy_orchestrator()
    res = _lazy_research()

    futures = {}
    with ThreadPoolExecutor(max_workers=_MAX_WORKERS) as executor:
        # Wikipedia + Clearbit logo
        if api:
            futures["company_info"] = executor.submit(
                _safe_call, api.fetch_company_info, company_name
            )
            futures["metadata"] = executor.submit(
                _safe_call, api.fetch_company_metadata, company_name
            )
            futures["sec"] = executor.submit(
                _safe_call, api.fetch_sec_company_data, company_name
            )

        # Employer brand from orchestrator
        if orch:
            futures["employer_brand"] = executor.submit(
                _safe_call, orch.enrich_employer_brand, company_name
            )

        # Company intelligence from research
        if res:
            futures["intelligence"] = executor.submit(
                _safe_call, res.get_company_intelligence, company_name
            )

        # Collect results
        for key, future in futures.items():
            try:
                result = future.result(timeout=15)
                if not result:
                    continue

                if key == "company_info":
                    if result.get("description"):
                        # Validate entity: check the first sentence
                        # mentions the company name to catch wrong-entity
                        # returns (e.g. "Square Co." video game dev for
                        # a localization company).
                        _desc = result["description"]
                        _first = (
                            _desc.split(".")[0].lower()
                            if "." in _desc
                            else _desc.lower()
                        )
                        _name_tokens = {
                            t
                            for t in re.split(r"[\s\-_.,/&]+", company_name.lower())
                            if t
                            and len(t) > 2
                            and t
                            not in {
                                "inc",
                                "llc",
                                "ltd",
                                "corp",
                                "the",
                                "and",
                                "company",
                                "co",
                            }
                        }
                        _has_match = (
                            any(tok in _first for tok in _name_tokens)
                            if _name_tokens
                            else True
                        )
                        if _has_match:
                            profile["description"] = _desc
                        else:
                            logger.warning(
                                "Entity mismatch in competitive_intel for '%s': "
                                "first sentence '%s' does not mention company name tokens %s",
                                company_name,
                                _first[:120],
                                _name_tokens,
                            )
                            profile["_entity_mismatch"] = True
                    if result.get("logo_url"):
                        profile["logo_url"] = result["logo_url"]
                    profile["sources"].append("Wikipedia/Clearbit")

                elif key == "metadata":
                    if result.get("domain"):
                        profile["domain"] = result["domain"]
                    if result.get("logo"):
                        profile["logo_url"] = profile["logo_url"] or result["logo"]
                    if result.get("name"):
                        # Use official name casing from Clearbit
                        profile["name"] = result["name"]
                    profile["sources"].append("Clearbit Autocomplete")

                elif key == "sec":
                    if isinstance(result, dict):
                        if result.get("ticker"):
                            profile["stock_ticker"] = result["ticker"]
                            profile["is_public"] = True
                        if result.get("sic_description"):
                            profile["industry"] = (
                                profile["industry"] or result["sic_description"]
                            )
                        if result.get("state_of_incorporation"):
                            profile["headquarters"] = (
                                profile["headquarters"]
                                or result["state_of_incorporation"]
                            )
                        profile["sources"].append("SEC EDGAR")

                elif key == "employer_brand":
                    if result.get("glassdoor_rating"):
                        profile["glassdoor_rating"] = result["glassdoor_rating"]
                    if result.get("employer_brand_strength"):
                        profile["employer_brand_strength"] = result[
                            "employer_brand_strength"
                        ]
                    if result.get("company_size"):
                        profile["employee_count"] = result["company_size"]
                    profile["sources"].append("Employer Brand DB")

                elif key == "intelligence":
                    if result.get("matched"):
                        if result.get("industry"):
                            profile["industry"] = (
                                profile["industry"] or result["industry"]
                            )
                        if result.get("employee_count"):
                            profile["employee_count"] = (
                                profile["employee_count"] or result["employee_count"]
                            )
                        if result.get("founded"):
                            profile["founded"] = result["founded"]
                        if result.get("headquarters"):
                            profile["headquarters"] = (
                                profile["headquarters"] or result["headquarters"]
                            )
                        profile["sources"].append("Company Intelligence DB")

            except Exception as exc:
                logger.warning("Failed to fetch %s for %s: %s", key, company_name, exc)

    # ── Fallback: fill gaps from well-known companies database ──
    _has_meaningful_data = bool(
        profile.get("description")
        or profile.get("employee_count")
        or profile.get("founded")
        or profile.get("headquarters")
    )
    if not _has_meaningful_data:
        fallback = _lookup_well_known_company(company_name)
        if fallback:
            logger.info("Using well-known company fallback for '%s'", company_name)
            for field in (
                "description",
                "industry",
                "employee_count",
                "founded",
                "headquarters",
                "domain",
                "stock_ticker",
                "is_public",
                "name",
            ):
                if fallback.get(field) and not profile.get(field):
                    profile[field] = fallback[field]
            # is_public needs special handling since False is falsy
            if fallback.get("is_public") and not profile.get("is_public"):
                profile["is_public"] = True
            if "Well-Known DB" not in profile.get("sources", []):
                profile["sources"].append("Well-Known DB")
    else:
        # Even with partial API data, fill remaining empty fields from fallback
        fallback = _lookup_well_known_company(company_name)
        if fallback:
            for field in (
                "description",
                "industry",
                "employee_count",
                "founded",
                "headquarters",
                "domain",
                "stock_ticker",
            ):
                if fallback.get(field) and not profile.get(field):
                    profile[field] = fallback[field]
            if fallback.get("is_public") and not profile.get("is_public"):
                profile["is_public"] = True
                if fallback.get("stock_ticker"):
                    profile["stock_ticker"] = (
                        profile.get("stock_ticker") or fallback["stock_ticker"]
                    )

    return profile


# ═══════════════════════════════════════════════════════════════════════════════
# 2. COMPETITOR COMPARISON
# ═══════════════════════════════════════════════════════════════════════════════


def analyze_competitors(
    company_name: str,
    competitor_names: List[str],
) -> Dict[str, Any]:
    """Analyze company and all competitors in parallel.

    Returns dict with:
      - company: profile dict for the primary company
      - competitors: list of profile dicts for each competitor
      - comparison_matrix: summary comparison table data
    """
    all_names = [company_name] + [c.strip() for c in competitor_names if c.strip()]

    # Fetch all profiles concurrently
    profiles: Dict[str, Dict[str, Any]] = {}
    with ThreadPoolExecutor(max_workers=_MAX_WORKERS) as executor:
        future_map = {
            executor.submit(analyze_company, name): name for name in all_names
        }
        for future in as_completed(future_map):
            name = future_map[future]
            try:
                profiles[name] = future.result(timeout=30)
            except Exception as exc:
                logger.warning("Failed to analyze %s: %s", name, exc)
                profiles[name] = {
                    "name": name,
                    "error": "Analysis failed for this company",
                }

    company_profile = profiles.get(company_name, {"name": company_name})
    competitor_profiles = [
        profiles.get(c, {"name": c}) for c in competitor_names if c.strip()
    ]

    # Build comparison matrix
    comparison = _build_comparison_matrix(company_profile, competitor_profiles)

    return {
        "company": company_profile,
        "competitors": competitor_profiles,
        "comparison_matrix": comparison,
    }


def _build_comparison_matrix(
    company: Dict[str, Any],
    competitors: List[Dict[str, Any]],
) -> List[Dict[str, Any]]:
    """Build a structured comparison table across all entities."""
    matrix = []

    def _row(profile: Dict[str, Any], is_primary: bool = False) -> Dict[str, Any]:
        emp = profile.get("employee_count")
        emp_str = _format_employee_count(emp) if emp else "N/A"
        return {
            "name": profile.get("name", "Unknown"),
            "is_primary": is_primary,
            "industry": profile.get("industry", "N/A"),
            "employee_count": emp_str,
            "employee_count_raw": emp if isinstance(emp, (int, float)) else 0,
            "is_public": profile.get("is_public", False),
            "stock_ticker": profile.get("stock_ticker") or "",
            "headquarters": profile.get("headquarters", "N/A"),
            "founded": profile.get("founded", "N/A"),
            "glassdoor_rating": profile.get("glassdoor_rating"),
            "employer_brand_strength": profile.get("employer_brand_strength") or "",
            "logo_url": profile.get("logo_url") or "",
        }

    matrix.append(_row(company, is_primary=True))
    for comp in competitors:
        matrix.append(_row(comp))

    return matrix


# ═══════════════════════════════════════════════════════════════════════════════
# 3. HIRING ACTIVITY COMPARISON
# ═══════════════════════════════════════════════════════════════════════════════


def compare_hiring_activity(
    company_name: str,
    competitors: List[str],
    industry: str = "general_entry_level",
    roles: Optional[List[str]] = None,
) -> Dict[str, Any]:
    """Compare hiring velocity and difficulty across company and competitors.

    Uses research.get_competitors and data_orchestrator.enrich_competitive
    for hiring intelligence. Returns hiring difficulty scores, competitive
    threat levels, and recommended channels per competitor.
    """
    result: Dict[str, Any] = {
        "company": company_name,
        "industry": industry,
        "industry_label": INDUSTRY_LABEL_MAP.get(industry, industry),
        "competitors_hiring": [],
        "industry_competitors": [],
        "hiring_difficulty": "moderate",
        "sources": [],
    }

    res = _lazy_research()
    orch = _lazy_orchestrator()

    # Get industry competitors from research.py
    if res:
        try:
            industry_comps = res.get_competitors(industry, [], company_name)
            if industry_comps:
                result["industry_competitors"] = industry_comps
                result["sources"].append("Industry Competitor DB")
        except Exception as exc:
            logger.warning("get_competitors failed: %s", exc)

    # Get competitive enrichment from orchestrator
    if orch:
        try:
            comp_data = orch.enrich_competitive(company_name, industry)
            if comp_data:
                if comp_data.get("competitors"):
                    # Merge industry competitor data
                    result["industry_competitors"] = (
                        result["industry_competitors"] or comp_data["competitors"]
                    )
                if comp_data.get("employer_brand"):
                    brand = comp_data["employer_brand"]
                    if brand.get("hiring_channels"):
                        result["recommended_channels"] = brand["hiring_channels"]
                result["sources"].append("Competitive Enrichment")
        except Exception as exc:
            logger.warning("enrich_competitive failed: %s", exc)

    # Get per-competitor intelligence
    if res and competitors:
        try:
            comp_intel = res.get_client_competitor_intelligence(competitors, industry)
            if comp_intel:
                result["competitors_hiring"] = comp_intel
                result["sources"].append("Competitor Intelligence")
        except Exception as exc:
            logger.warning("get_client_competitor_intelligence failed: %s", exc)

    # Determine overall hiring difficulty based on available data
    result["hiring_difficulty"] = _assess_hiring_difficulty(
        result.get("industry_competitors") or [],
        result.get("competitors_hiring") or [],
        industry,
    )

    return result


def _assess_hiring_difficulty(
    industry_comps: List[Any],
    competitor_intel: List[Any],
    industry: str,
) -> str:
    """Heuristic hiring difficulty rating based on competitor landscape."""
    high_threat_industries = {
        "tech_engineering",
        "healthcare_medical",
        "aerospace_defense",
        "pharma_biotech",
        "finance_banking",
    }
    if industry in high_threat_industries:
        return "high"

    # Count high-threat competitors
    high_threats = 0
    for comp in industry_comps:
        if isinstance(comp, dict) and comp.get("threat") in ("high", "critical"):
            high_threats += 1

    if high_threats >= 3:
        return "high"
    elif high_threats >= 1:
        return "moderate"
    return "low"


# ═══════════════════════════════════════════════════════════════════════════════
# 4. AD BENCHMARK COMPARISON
# ═══════════════════════════════════════════════════════════════════════════════


def compare_ad_benchmarks(
    industry: str = "general_entry_level",
    roles: Optional[List[str]] = None,
    collar_type: str = "mixed",
) -> Dict[str, Any]:
    """Get CPC/CPA benchmarks across all ad platforms for the industry.

    Returns per-platform benchmark data including CPC, CPA, CTR,
    trend direction, and confidence intervals.
    """
    result: Dict[str, Any] = {
        "industry": industry,
        "industry_label": INDUSTRY_LABEL_MAP.get(industry, industry),
        "collar_type": collar_type,
        "platforms": {},
        "sources": [],
    }

    te = _lazy_trends()
    if te:
        try:
            benchmarks = te.get_all_platform_benchmarks(
                industry=industry,
                collar_type=collar_type,
            )
            if benchmarks:
                for plat_key, data in benchmarks.items():
                    result["platforms"][plat_key] = {
                        "name": PLATFORM_DISPLAY.get(plat_key, plat_key),
                        "cpc": round(data.get("value") or 0, 2),
                        "cpa": round(data.get("cpa_value") or 0, 2),
                        "ctr": (
                            round((data.get("ctr_value") or 0) * 100, 2)
                            if data.get("ctr_value")
                            else None
                        ),
                        "trend_direction": data.get("trend_direction", "stable"),
                        "trend_pct_yoy": data.get("trend_pct_yoy"),
                        "confidence": data.get("data_confidence", 0.5),
                        "seasonal_factor": data.get("seasonal_factor", 1.0),
                    }
                result["sources"].append("Trend Engine (Appcast/WordStream/SHRM)")
        except Exception as exc:
            logger.warning("get_all_platform_benchmarks failed: %s", exc)

    # Fallback benchmarks when trend_engine unavailable
    if not result["platforms"]:
        result["platforms"] = _fallback_benchmarks(industry)
        result["sources"].append("Fallback Benchmarks")

    return result


def _fallback_benchmarks(industry: str) -> Dict[str, Dict[str, Any]]:
    """Provide reasonable fallback benchmarks when trend_engine is unavailable.

    Prefers benchmark_registry (single source of truth) when available,
    otherwise falls back to hardcoded values for resilience.
    """
    # -- Prefer unified benchmark_registry --
    if _HAS_BENCHMARK_REGISTRY:
        _platforms = [
            "google_search",
            "meta_facebook",
            "meta_instagram",
            "linkedin",
            "indeed",
            "programmatic",
        ]
        result: Dict[str, Dict[str, Any]] = {}
        for plat_key in _platforms:
            bench = get_channel_benchmark(plat_key, industry)
            result[plat_key] = {
                "name": PLATFORM_DISPLAY.get(plat_key, plat_key),
                "cpc": bench.get("cpc_adjusted") or bench.get("cpc", 1.0),
                "cpa": bench.get("cpa_adjusted") or bench.get("cpa", 30.0),
                "ctr": bench.get("ctr", 0.025),
                "trend_direction": "stable",
                "trend_pct_yoy": None,
                "confidence": (
                    0.6 if bench.get("data_source") == "live_firecrawl" else 0.4
                ),
                "seasonal_factor": 1.0,
            }
        return result

    # -- Hardcoded fallback (kept for resilience if registry unavailable) --
    industry_multipliers = {
        "healthcare_medical": 1.2,
        "tech_engineering": 1.4,
        "finance_banking": 1.3,
        "legal_services": 1.5,
        "aerospace_defense": 1.3,
        "pharma_biotech": 1.35,
        "blue_collar_trades": 0.7,
        "general_entry_level": 0.6,
        "retail_consumer": 0.65,
        "hospitality_travel": 0.7,
        "food_beverage": 0.6,
        "logistics_supply_chain": 0.8,
        "construction_real_estate": 0.85,
        "education": 0.9,
    }
    mult = industry_multipliers.get(industry, 1.0)

    base = {
        "google_search": {"cpc": 2.69, "cpa": 45.00, "ctr": 3.2},
        "meta_facebook": {"cpc": 1.72, "cpa": 30.00, "ctr": 1.1},
        "meta_instagram": {"cpc": 1.50, "cpa": 35.00, "ctr": 0.9},
        "linkedin": {
            "cpc": 5.26,
            "cpa": 45.00,
            "ctr": 0.5,
        },  # Sponsored Jobs CPA $30-$90, US avg $45 (2026-04-07)
        "indeed": {"cpc": 0.50, "cpa": 25.00, "ctr": 4.5},
        "programmatic": {"cpc": 0.63, "cpa": 22.00, "ctr": 2.8},
    }

    result = {}
    for plat_key, vals in base.items():
        result[plat_key] = {
            "name": PLATFORM_DISPLAY.get(plat_key, plat_key),
            "cpc": round(vals["cpc"] * mult, 2),
            "cpa": round(vals["cpa"] * mult, 2),
            "ctr": round(vals["ctr"], 2),
            "trend_direction": "stable",
            "trend_pct_yoy": None,
            "confidence": 0.4,
            "seasonal_factor": 1.0,
        }
    return result


# ═══════════════════════════════════════════════════════════════════════════════
# 5. MARKET TRENDS (Google Trends)
# ═══════════════════════════════════════════════════════════════════════════════


def get_market_trends(
    company_name: str,
    competitors: List[str],
) -> Dict[str, Any]:
    """Compare Google Trends search interest for company vs competitors.

    Uses api_enrichment.fetch_search_trends to get relative search volume.
    Returns per-company interest scores with trend direction.
    """
    result: Dict[str, Any] = {
        "companies": {},
        "source": "",
        "max_interest": 0,
    }

    all_names = [company_name] + [c.strip() for c in competitors if c.strip()]
    # Google Trends supports max 5 keywords at once
    keywords = all_names[:5]

    api = _lazy_api()
    if api:
        try:
            trends = api.fetch_search_trends(keywords)
            if trends and len(trends) > 1:
                result["source"] = trends.get("source", "Google Trends")
                max_interest = 0
                for kw in keywords:
                    if kw in trends:
                        entry = trends[kw]
                        avg = entry.get("avg_interest", 50)
                        latest = entry.get("latest_interest", avg)
                        trend_dir = entry.get("trend", "stable")
                        result["companies"][kw] = {
                            "avg_interest": avg,
                            "latest_interest": latest,
                            "trend": trend_dir,
                            "is_primary": kw == company_name,
                        }
                        max_interest = max(max_interest, avg, latest)
                result["max_interest"] = max_interest
        except Exception as exc:
            logger.warning("fetch_search_trends failed: %s", exc)

    # Fallback: assign synthetic relative interest based on name length hash
    if not result["companies"]:
        result["source"] = "Estimated (API unavailable)"
        for i, kw in enumerate(keywords):
            base = 60 if i == 0 else max(20, 80 - (i * 15))
            # Simple deterministic variation based on company name
            variation = (sum(ord(c) for c in kw) % 30) - 15
            interest = max(10, min(100, base + variation))
            result["companies"][kw] = {
                "avg_interest": interest,
                "latest_interest": max(10, interest + (variation // 3)),
                "trend": (
                    "rising"
                    if variation > 5
                    else ("declining" if variation < -5 else "stable")
                ),
                "is_primary": kw == company_name,
            }
        result["max_interest"] = (
            max(v["avg_interest"] for v in result["companies"].values())
            if result["companies"]
            else 100
        )

    return result


# ═══════════════════════════════════════════════════════════════════════════════
# 6. GENERATE COMPETITIVE BRIEF
# ═══════════════════════════════════════════════════════════════════════════════


def generate_competitive_brief(analysis_results: Dict[str, Any]) -> Dict[str, Any]:
    """Compile all analysis data into a structured report dict.

    Generates strategic recommendations based on the competitive landscape.
    """
    brief: Dict[str, Any] = {
        "generated_at": datetime.utcnow().isoformat() + "Z",
        "company": analysis_results.get("competitor_analysis", {}).get("company", {}),
        "competitors": analysis_results.get("competitor_analysis", {}).get(
            "competitors", []
        ),
        "comparison_matrix": analysis_results.get("competitor_analysis", {}).get(
            "comparison_matrix", []
        ),
        "hiring_activity": analysis_results.get("hiring_activity", {}),
        "ad_benchmarks": analysis_results.get("ad_benchmarks", {}),
        "market_trends": analysis_results.get("market_trends", {}),
        "recommendations": [],
    }

    # Generate strategic recommendations
    brief["recommendations"] = _generate_recommendations(brief)

    # Generate AI competitive narrative
    brief["competitive_narrative"] = _generate_competitive_narrative(
        {
            "company": brief.get("company", {}),
            "competitors": brief.get("competitors") or [][:5],  # Cap for token budget
            "hiring_activity": brief.get("hiring_activity", {}),
            "ad_benchmarks": brief.get("ad_benchmarks", {}),
        }
    )

    return brief


def _generate_recommendations(brief: Dict[str, Any]) -> List[Dict[str, str]]:
    """Generate actionable recommendations from the analysis data."""
    recs: List[Dict[str, str]] = []

    # 1. Channel recommendation based on CPC/CPA
    benchmarks = brief.get("ad_benchmarks", {}).get("platforms", {})
    if benchmarks:
        # Find lowest CPA platform
        lowest_cpa_plat = None
        lowest_cpa = float("inf")
        highest_ctr_plat = None
        highest_ctr = 0

        for plat_key, data in benchmarks.items():
            cpa = data.get("cpa", float("inf"))
            ctr = data.get("ctr") or 0 or 0
            if cpa < lowest_cpa:
                lowest_cpa = cpa
                lowest_cpa_plat = data.get("name", plat_key)
            if ctr > highest_ctr:
                highest_ctr = ctr
                highest_ctr_plat = data.get("name", plat_key)

        if lowest_cpa_plat:
            recs.append(
                {
                    "title": "Optimize for Cost Efficiency",
                    "description": (
                        f"{lowest_cpa_plat} offers the lowest cost-per-application "
                        f"(${lowest_cpa:.2f}) in your industry. Allocate 40-50% of "
                        f"budget here for maximum applicant volume."
                    ),
                    "priority": "high",
                    "icon": "dollar",
                }
            )

        if highest_ctr_plat and highest_ctr_plat != lowest_cpa_plat:
            recs.append(
                {
                    "title": "Maximize Engagement",
                    "description": (
                        f"{highest_ctr_plat} shows the highest click-through rate "
                        f"({highest_ctr:.1f}%). Use compelling creative and "
                        f"employer branding content on this platform."
                    ),
                    "priority": "medium",
                    "icon": "chart",
                }
            )

    # 2. Competitive positioning
    hiring = brief.get("hiring_activity", {})
    difficulty = hiring.get("hiring_difficulty", "moderate")
    if difficulty == "high":
        recs.append(
            {
                "title": "Strengthen Employer Brand",
                "description": (
                    "Your industry has high hiring competition. Invest in employer "
                    "branding, Glassdoor profile optimization, and employee advocacy "
                    "programs to differentiate from competitors."
                ),
                "priority": "high",
                "icon": "shield",
            }
        )
    elif difficulty == "low":
        recs.append(
            {
                "title": "Capitalize on Low Competition",
                "description": (
                    "Hiring competition is relatively low in your space. Focus on "
                    "volume-driven channels and programmatic distribution to fill "
                    "positions quickly at favorable CPAs."
                ),
                "priority": "medium",
                "icon": "rocket",
            }
        )

    # 3. Search interest trend
    trends = brief.get("market_trends", {})
    companies = trends.get("companies", {})
    company_profile = brief.get("company", {})
    company_name = company_profile.get("name") or ""
    if company_name in companies:
        trend = companies[company_name].get("trend", "stable")
        if trend == "declining":
            recs.append(
                {
                    "title": "Boost Brand Awareness",
                    "description": (
                        "Search interest for your company is declining relative to "
                        "competitors. Consider increasing investment in awareness "
                        "campaigns, social media presence, and content marketing."
                    ),
                    "priority": "high",
                    "icon": "megaphone",
                }
            )
        elif trend == "rising":
            recs.append(
                {
                    "title": "Leverage Rising Interest",
                    "description": (
                        "Search interest for your company is trending upward. "
                        "Capitalize on this momentum with targeted recruitment "
                        "campaigns and career page optimization."
                    ),
                    "priority": "medium",
                    "icon": "trending",
                }
            )

    # 4. Diversification recommendation
    if len(benchmarks) >= 4:
        recs.append(
            {
                "title": "Diversify Channel Mix",
                "description": (
                    "Avoid over-reliance on a single platform. A balanced mix of "
                    "job boards (Indeed), social (LinkedIn, Meta), search (Google), "
                    "and programmatic channels reduces risk and reaches passive candidates."
                ),
                "priority": "medium",
                "icon": "grid",
            }
        )

    # 5. Competitor count check
    comp_count = len(brief.get("competitors") or [])
    if comp_count >= 3:
        recs.append(
            {
                "title": "Monitor Competitor Moves",
                "description": (
                    f"With {comp_count} key competitors tracked, set up regular "
                    f"monitoring of their career pages, Glassdoor reviews, and "
                    f"job posting volumes to stay ahead of hiring surges."
                ),
                "priority": "low",
                "icon": "eye",
            }
        )

    return recs[:5]  # Cap at 5 recommendations


# ═══════════════════════════════════════════════════════════════════════════════
# 7. EXCEL REPORT GENERATION
# ═══════════════════════════════════════════════════════════════════════════════


def generate_competitive_excel(
    brief: Dict[str, Any],
    company_name: str,
) -> bytes:
    """Generate a competitive intelligence Excel workbook.

    Sheets:
      1. Company Overview - profile, key facts
      2. Competitor Comparison - side-by-side matrix
      3. Market Position - benchmarks, trends
      4. Recommendations - strategic actions

    Uses Sapphire Blue palette (consistent with excel_v2.py).
    """
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    # Design tokens -- Sapphire Blue palette
    NAVY = "0F172A"
    SAPPHIRE = "2563EB"
    BLUE_LIGHT = "DBEAFE"
    BLUE_PALE = "EFF6FF"
    GREEN = "16A34A"
    WARM_GRAY = "E7E5E4"
    OFF_WHITE = "F5F5F4"
    MUTED = "78716C"

    hdr_fill = PatternFill(start_color=NAVY, end_color=NAVY, fill_type="solid")
    hdr_font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
    accent_fill = PatternFill(
        start_color=BLUE_PALE, end_color=BLUE_PALE, fill_type="solid"
    )
    data_font = Font(name="Calibri", size=10, color="1E293B")
    bold_font = Font(name="Calibri", size=10, bold=True, color="1E293B")
    title_font = Font(name="Calibri", size=14, bold=True, color=NAVY)
    subtitle_font = Font(name="Calibri", size=11, color=MUTED)
    thin_border = Border(
        bottom=Side(style="thin", color=WARM_GRAY),
    )
    wrap_align = Alignment(wrap_text=True, vertical="top")
    center_align = Alignment(horizontal="center", vertical="center")

    wb = Workbook()

    # ── Sheet 1: Company Overview ──
    ws1 = wb.active
    ws1.title = "Company Overview"
    ws1.sheet_properties.tabColor = SAPPHIRE

    company = brief.get("company", {})
    ws1.merge_cells("B2:F2")
    ws1["B2"] = f"Competitive Intelligence: {company_name}"
    ws1["B2"].font = title_font

    ws1.merge_cells("B3:F3")
    ws1["B3"] = (
        f"Generated {brief.get('generated_at') or ''[:10]} | Powered by Nova AI Suite"
    )
    ws1["B3"].font = subtitle_font

    # Company profile table
    row = 5
    profile_fields = [
        ("Company Name", company.get("name", company_name)),
        ("Industry", company.get("industry", "N/A")),
        ("Employees", _format_employee_count(company.get("employee_count"))),
        ("Founded", str(company.get("founded", "N/A"))),
        ("Headquarters", company.get("headquarters", "N/A")),
        ("Stock Ticker", company.get("stock_ticker", "N/A") or "Private"),
        ("Glassdoor Rating", str(company.get("glassdoor_rating", "N/A") or "N/A")),
        ("Employer Brand", company.get("employer_brand_strength", "N/A") or "N/A"),
        ("Domain", company.get("domain", "N/A") or "N/A"),
    ]

    for label, value in profile_fields:
        ws1[f"B{row}"] = label
        ws1[f"B{row}"].font = bold_font
        ws1[f"B{row}"].fill = accent_fill
        ws1.merge_cells(f"C{row}:F{row}")
        ws1[f"C{row}"] = str(value)
        ws1[f"C{row}"].font = data_font
        ws1[f"C{row}"].border = thin_border
        row += 1

    # Description
    row += 1
    ws1[f"B{row}"] = "Description"
    ws1[f"B{row}"].font = bold_font
    row += 1
    desc = company.get("description", "No description available.")
    ws1.merge_cells(f"B{row}:F{row + 2}")
    ws1[f"B{row}"] = desc[:500] if desc else "No description available."
    ws1[f"B{row}"].font = data_font
    ws1[f"B{row}"].alignment = wrap_align

    ws1.column_dimensions["A"].width = 3
    ws1.column_dimensions["B"].width = 20
    for col in "CDEF":
        ws1.column_dimensions[col].width = 18

    # ── Sheet 2: Competitor Comparison ──
    ws2 = wb.create_sheet("Competitor Comparison")
    ws2.sheet_properties.tabColor = SAPPHIRE

    ws2.merge_cells("B2:H2")
    ws2["B2"] = "Competitor Comparison Matrix"
    ws2["B2"].font = title_font

    matrix = brief.get("comparison_matrix") or []
    if matrix:
        headers = [
            "Company",
            "Industry",
            "Employees",
            "Public/Private",
            "Ticker",
            "HQ",
            "Glassdoor",
            "Brand Strength",
        ]
        for ci, h in enumerate(headers, start=2):
            cell = ws2.cell(row=4, column=ci, value=h)
            cell.font = hdr_font
            cell.fill = hdr_fill
            cell.alignment = center_align

        for ri, entry in enumerate(matrix, start=5):
            ws2.cell(row=ri, column=2, value=entry.get("name") or "").font = (
                bold_font if entry.get("is_primary") else data_font
            )
            ws2.cell(row=ri, column=3, value=entry.get("industry", "N/A")).font = (
                data_font
            )
            ws2.cell(
                row=ri, column=4, value=entry.get("employee_count", "N/A")
            ).font = data_font
            ws2.cell(
                row=ri,
                column=5,
                value="Public" if entry.get("is_public") else "Private",
            ).font = data_font
            ws2.cell(
                row=ri, column=6, value=entry.get("stock_ticker") or "" or "-"
            ).font = data_font
            ws2.cell(row=ri, column=7, value=entry.get("headquarters", "N/A")).font = (
                data_font
            )
            ws2.cell(
                row=ri, column=8, value=str(entry.get("glassdoor_rating", "-") or "-")
            ).font = data_font
            ws2.cell(
                row=ri, column=9, value=entry.get("employer_brand_strength", "-") or "-"
            ).font = data_font

            # Highlight primary company row
            if entry.get("is_primary"):
                for cc in range(2, 10):
                    ws2.cell(row=ri, column=cc).fill = accent_fill

    ws2.column_dimensions["A"].width = 3
    for col_idx in range(2, 10):
        ws2.column_dimensions[get_column_letter(col_idx)].width = 18

    # ── Sheet 3: Market Position ──
    ws3 = wb.create_sheet("Market Position")
    ws3.sheet_properties.tabColor = SAPPHIRE

    ws3.merge_cells("B2:H2")
    ws3["B2"] = "Ad Platform Benchmarks & Market Trends"
    ws3["B2"].font = title_font

    # Ad benchmarks table
    benchmarks = brief.get("ad_benchmarks", {}).get("platforms", {})
    if benchmarks:
        ws3["B4"] = "Platform Benchmarks"
        ws3["B4"].font = bold_font

        bench_headers = ["Platform", "CPC", "CPA", "CTR (%)", "Trend", "YoY Change"]
        for ci, h in enumerate(bench_headers, start=2):
            cell = ws3.cell(row=5, column=ci, value=h)
            cell.font = hdr_font
            cell.fill = hdr_fill
            cell.alignment = center_align

        row = 6
        for plat_key, data in benchmarks.items():
            ws3.cell(row=row, column=2, value=data.get("name", plat_key)).font = (
                bold_font
            )
            ws3.cell(row=row, column=3, value=f"${data.get('cpc') or 0:.2f}").font = (
                data_font
            )
            ws3.cell(row=row, column=4, value=f"${data.get('cpa') or 0:.2f}").font = (
                data_font
            )
            ctr_val = data.get("ctr")
            ws3.cell(
                row=row, column=5, value=f"{ctr_val:.1f}%" if ctr_val else "N/A"
            ).font = data_font
            ws3.cell(
                row=row, column=6, value=data.get("trend_direction", "stable").title()
            ).font = data_font
            yoy = data.get("trend_pct_yoy")
            ws3.cell(row=row, column=7, value=f"{yoy:+.1f}%" if yoy else "N/A").font = (
                data_font
            )
            row += 1

    # Market trends
    trends = brief.get("market_trends", {}).get("companies", {})
    if trends:
        row += 2
        ws3[f"B{row}"] = "Search Interest Comparison"
        ws3[f"B{row}"].font = bold_font
        row += 1

        trend_headers = ["Company", "Avg Interest", "Latest Interest", "Trend"]
        for ci, h in enumerate(trend_headers, start=2):
            cell = ws3.cell(row=row, column=ci, value=h)
            cell.font = hdr_font
            cell.fill = hdr_fill
            cell.alignment = center_align
        row += 1

        for name, data in trends.items():
            ws3.cell(row=row, column=2, value=name).font = (
                bold_font if data.get("is_primary") else data_font
            )
            ws3.cell(row=row, column=3, value=data.get("avg_interest", "N/A")).font = (
                data_font
            )
            ws3.cell(
                row=row, column=4, value=data.get("latest_interest", "N/A")
            ).font = data_font
            ws3.cell(
                row=row, column=5, value=data.get("trend", "stable").title()
            ).font = data_font
            if data.get("is_primary"):
                for cc in range(2, 6):
                    ws3.cell(row=row, column=cc).fill = accent_fill
            row += 1

    ws3.column_dimensions["A"].width = 3
    for col_idx in range(2, 9):
        ws3.column_dimensions[get_column_letter(col_idx)].width = 18

    # ── Sheet 4: Recommendations ──
    ws4 = wb.create_sheet("Recommendations")
    ws4.sheet_properties.tabColor = GREEN

    ws4.merge_cells("B2:F2")
    ws4["B2"] = "Strategic Recommendations"
    ws4["B2"].font = title_font

    recs = brief.get("recommendations") or []
    row = 4
    for i, rec in enumerate(recs, start=1):
        priority = rec.get("priority", "medium")
        pri_colors = {"high": "DC2626", "medium": "F59E0B", "low": "16A34A"}
        pri_fill = PatternFill(
            start_color=pri_colors.get(priority, "F59E0B"),
            end_color=pri_colors.get(priority, "F59E0B"),
            fill_type="solid",
        )

        ws4[f"B{row}"] = f"#{i}"
        ws4[f"B{row}"].font = Font(name="Calibri", size=10, bold=True, color="FFFFFF")
        ws4[f"B{row}"].fill = pri_fill
        ws4[f"B{row}"].alignment = center_align

        ws4[f"C{row}"] = rec.get("title") or ""
        ws4[f"C{row}"].font = bold_font

        ws4[f"D{row}"] = priority.upper()
        ws4[f"D{row}"].font = Font(
            name="Calibri", size=9, bold=True, color=pri_colors.get(priority, "F59E0B")
        )

        row += 1
        ws4.merge_cells(f"C{row}:F{row}")
        ws4[f"C{row}"] = rec.get("description") or ""
        ws4[f"C{row}"].font = data_font
        ws4[f"C{row}"].alignment = wrap_align
        row += 2

    ws4.column_dimensions["A"].width = 3
    ws4.column_dimensions["B"].width = 5
    ws4.column_dimensions["C"].width = 30
    ws4.column_dimensions["D"].width = 12
    for col in "EF":
        ws4.column_dimensions[col].width = 20

    # Save
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ═══════════════════════════════════════════════════════════════════════════════
# 8. POWERPOINT REPORT GENERATION
# ═══════════════════════════════════════════════════════════════════════════════


def generate_competitive_ppt(
    brief: Dict[str, Any],
    company_name: str,
) -> bytes:
    """Generate a competitive intelligence PowerPoint presentation.

    Slides:
      1. Title slide (company name, date, branding)
      2. Company Overview (profile card)
      3. Competitive Landscape (comparison table)
      4. Market Trends (search interest bars)
      5. Ad Platform Benchmarks (CPC/CPA comparison)
      6. Recommendations (strategic actions)

    Uses brand identity: Port Gore navy, Blue Violet purple, Downy teal.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

    # Brand colors
    PORT_GORE = RGBColor(0x20, 0x20, 0x58)
    BLUE_VIOLET = RGBColor(0x5A, 0x54, 0xBD)
    DOWNY_TEAL = RGBColor(0x6B, 0xB3, 0xCD)
    TAPESTRY_PINK = RGBColor(0xB5, 0x66, 0x9C)
    RAW_SIENNA = RGBColor(0xCE, 0x90, 0x47)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF5)
    DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
    MUTED_TEXT = RGBColor(0x64, 0x74, 0x8B)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def _add_bg(slide, color=PORT_GORE):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_text_box(
        slide,
        left,
        top,
        width,
        height,
        text,
        font_size=12,
        bold=False,
        color=WHITE,
        alignment=PP_ALIGN.LEFT,
    ):
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = alignment
        return txBox

    def _add_shape(slide, left, top, width, height, color):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        shape.shadow.inherit = False
        return shape

    # ── Slide 1: Title ──
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_bg(slide1)

    # Accent bar
    bar = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(3.0),
        Inches(13.333),
        Inches(0.06),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE_VIOLET
    bar.line.fill.background()

    _add_text_box(
        slide1,
        1,
        1.5,
        11,
        1.5,
        f"Competitive Intelligence Report",
        font_size=36,
        bold=True,
        color=WHITE,
    )
    _add_text_box(
        slide1,
        1,
        3.3,
        11,
        0.8,
        company_name,
        font_size=24,
        bold=False,
        color=DOWNY_TEAL,
    )
    _add_text_box(
        slide1,
        1,
        4.5,
        11,
        0.5,
        f"Generated {brief.get('generated_at') or ''[:10]} | Powered by Nova AI Suite",
        font_size=12,
        color=MUTED_TEXT,
    )

    # ── Slide 2: Company Overview ──
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide2, RGBColor(0x0F, 0x0F, 0x1E))

    _add_text_box(
        slide2,
        0.8,
        0.4,
        6,
        0.6,
        "Company Overview",
        font_size=24,
        bold=True,
        color=WHITE,
    )

    company = brief.get("company", {})
    profile_items = [
        ("Industry", company.get("industry", "N/A")),
        ("Employees", _format_employee_count(company.get("employee_count"))),
        ("Founded", str(company.get("founded", "N/A"))),
        ("Headquarters", company.get("headquarters", "N/A")),
        ("Status", "Public" if company.get("is_public") else "Private"),
        ("Ticker", company.get("stock_ticker") or "" or "N/A"),
        ("Glassdoor", str(company.get("glassdoor_rating", "N/A") or "N/A")),
    ]

    # Profile card background
    _add_shape(slide2, 0.8, 1.2, 5.5, 4.5, RGBColor(0x1A, 0x1A, 0x30))

    _add_text_box(
        slide2,
        1.1,
        1.4,
        5,
        0.5,
        company.get("name", company_name),
        font_size=20,
        bold=True,
        color=DOWNY_TEAL,
    )

    y = 2.1
    for label, value in profile_items:
        _add_text_box(
            slide2, 1.1, y, 2.2, 0.35, label, font_size=10, bold=True, color=MUTED_TEXT
        )
        _add_text_box(slide2, 3.3, y, 2.8, 0.35, str(value), font_size=10, color=WHITE)
        y += 0.35

    # Description card
    desc = company.get("description", "No description available.")
    _add_shape(slide2, 6.8, 1.2, 5.7, 4.5, RGBColor(0x1A, 0x1A, 0x30))
    _add_text_box(
        slide2, 7.0, 1.4, 5.3, 0.4, "About", font_size=14, bold=True, color=DOWNY_TEAL
    )
    _add_text_box(
        slide2,
        7.0,
        1.9,
        5.3,
        3.5,
        desc[:400] if desc else "No description available.",
        font_size=10,
        color=WHITE,
    )

    # ── Slide 3: Competitive Landscape ──
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide3, RGBColor(0x0F, 0x0F, 0x1E))

    _add_text_box(
        slide3,
        0.8,
        0.4,
        6,
        0.6,
        "Competitive Landscape",
        font_size=24,
        bold=True,
        color=WHITE,
    )

    matrix = brief.get("comparison_matrix") or []
    if matrix:
        # Table header
        col_headers = ["Company", "Industry", "Employees", "Status", "Glassdoor"]
        col_widths = [2.5, 2.8, 1.6, 1.2, 1.2]
        x_start = 0.8
        y_start = 1.3

        # Header row
        x = x_start
        for ci, (hdr, w) in enumerate(zip(col_headers, col_widths)):
            shape = _add_shape(slide3, x, y_start, w, 0.45, BLUE_VIOLET)
            _add_text_box(
                slide3,
                x + 0.1,
                y_start + 0.05,
                w - 0.2,
                0.35,
                hdr,
                font_size=10,
                bold=True,
                color=WHITE,
                alignment=PP_ALIGN.LEFT,
            )
            x += w

        # Data rows
        for ri, entry in enumerate(matrix):
            y = y_start + 0.5 + (ri * 0.45)
            row_color = (
                RGBColor(0x1A, 0x1A, 0x30)
                if ri % 2 == 0
                else RGBColor(0x15, 0x15, 0x28)
            )
            x = x_start
            values = [
                entry.get("name") or "",
                entry.get("industry", "N/A"),
                entry.get("employee_count", "N/A"),
                "Public" if entry.get("is_public") else "Private",
                str(entry.get("glassdoor_rating", "-") or "-"),
            ]
            for vi, (val, w) in enumerate(zip(values, col_widths)):
                _add_shape(slide3, x, y, w, 0.45, row_color)
                text_color = (
                    DOWNY_TEAL if vi == 0 and entry.get("is_primary") else WHITE
                )
                font_bold = vi == 0 and entry.get("is_primary")
                _add_text_box(
                    slide3,
                    x + 0.1,
                    y + 0.05,
                    w - 0.2,
                    0.35,
                    str(val),
                    font_size=9,
                    bold=font_bold,
                    color=text_color,
                    alignment=PP_ALIGN.LEFT,
                )
                x += w

    # ── Slide 4: Market Trends ──
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide4, RGBColor(0x0F, 0x0F, 0x1E))

    _add_text_box(
        slide4,
        0.8,
        0.4,
        6,
        0.6,
        "Market Trends & Search Interest",
        font_size=24,
        bold=True,
        color=WHITE,
    )

    trends = brief.get("market_trends", {}).get("companies", {})
    max_interest = brief.get("market_trends", {}).get("max_interest", 100) or 100

    if trends:
        y = 1.5
        bar_max_width = 8.0
        for name, data in trends.items():
            interest = data.get("avg_interest", 50)
            bar_width = max(0.3, (interest / max_interest) * bar_max_width)
            is_primary = data.get("is_primary", False)
            bar_color = DOWNY_TEAL if is_primary else BLUE_VIOLET
            trend_label = data.get("trend", "stable")

            _add_text_box(
                slide4,
                0.8,
                y,
                3,
                0.35,
                name,
                font_size=11,
                bold=is_primary,
                color=WHITE,
            )
            _add_shape(slide4, 4.0, y + 0.05, bar_width, 0.3, bar_color)
            _add_text_box(
                slide4,
                4.0 + bar_width + 0.2,
                y,
                1.5,
                0.35,
                f"{interest} ({trend_label})",
                font_size=9,
                color=MUTED_TEXT,
            )
            y += 0.6

    _add_text_box(
        slide4,
        0.8,
        6.5,
        6,
        0.3,
        f"Source: {brief.get('market_trends', {}).get('source', 'Google Trends')}",
        font_size=8,
        color=MUTED_TEXT,
    )

    # ── Slide 5: Ad Platform Benchmarks ──
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide5, RGBColor(0x0F, 0x0F, 0x1E))

    _add_text_box(
        slide5,
        0.8,
        0.4,
        6,
        0.6,
        "Ad Platform Benchmarks",
        font_size=24,
        bold=True,
        color=WHITE,
    )

    benchmarks = brief.get("ad_benchmarks", {}).get("platforms", {})
    if benchmarks:
        bench_cols = ["Platform", "CPC", "CPA", "CTR", "Trend"]
        bench_widths = [2.5, 1.5, 1.5, 1.5, 1.5]
        x_start = 0.8
        y_start = 1.3

        x = x_start
        for hdr, w in zip(bench_cols, bench_widths):
            _add_shape(slide5, x, y_start, w, 0.45, BLUE_VIOLET)
            _add_text_box(
                slide5,
                x + 0.1,
                y_start + 0.05,
                w - 0.2,
                0.35,
                hdr,
                font_size=10,
                bold=True,
                color=WHITE,
            )
            x += w

        for ri, (plat_key, data) in enumerate(benchmarks.items()):
            y = y_start + 0.5 + (ri * 0.45)
            row_color = (
                RGBColor(0x1A, 0x1A, 0x30)
                if ri % 2 == 0
                else RGBColor(0x15, 0x15, 0x28)
            )
            x = x_start

            ctr_val = data.get("ctr")
            values = [
                data.get("name", plat_key),
                f"${data.get('cpc') or 0:.2f}",
                f"${data.get('cpa') or 0:.2f}",
                f"{ctr_val:.1f}%" if ctr_val else "N/A",
                data.get("trend_direction", "stable").title(),
            ]
            for val, w in zip(values, bench_widths):
                _add_shape(slide5, x, y, w, 0.45, row_color)
                _add_text_box(
                    slide5,
                    x + 0.1,
                    y + 0.05,
                    w - 0.2,
                    0.35,
                    val,
                    font_size=9,
                    color=WHITE,
                )
                x += w

    industry_label = brief.get("ad_benchmarks", {}).get("industry_label") or ""
    if industry_label:
        _add_text_box(
            slide5,
            0.8,
            6.5,
            8,
            0.3,
            f"Industry: {industry_label} | Source: Appcast, WordStream, SHRM Benchmarks",
            font_size=8,
            color=MUTED_TEXT,
        )

    # ── Slide 6: Recommendations ──
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide6, RGBColor(0x0F, 0x0F, 0x1E))

    _add_text_box(
        slide6,
        0.8,
        0.4,
        6,
        0.6,
        "Strategic Recommendations",
        font_size=24,
        bold=True,
        color=WHITE,
    )

    recs = brief.get("recommendations") or []
    priority_colors = {
        "high": RGBColor(0xDC, 0x26, 0x26),
        "medium": RAW_SIENNA,
        "low": RGBColor(0x16, 0xA3, 0x4A),
    }

    y = 1.3
    card_width = 5.5
    for i, rec in enumerate(recs):
        # Two-column layout
        col = 0 if i % 2 == 0 else 1
        if i > 0 and col == 0:
            y += 1.8
        x = 0.8 if col == 0 else 7.0

        priority = rec.get("priority", "medium")
        _add_shape(slide6, x, y, card_width, 1.5, RGBColor(0x1A, 0x1A, 0x30))

        # Priority badge
        badge = _add_shape(
            slide6,
            x + 0.15,
            y + 0.15,
            0.7,
            0.3,
            priority_colors.get(priority, RAW_SIENNA),
        )

        _add_text_box(
            slide6,
            x + 0.15,
            y + 0.15,
            0.7,
            0.3,
            priority.upper(),
            font_size=7,
            bold=True,
            color=WHITE,
            alignment=PP_ALIGN.CENTER,
        )

        _add_text_box(
            slide6,
            x + 1.0,
            y + 0.1,
            card_width - 1.2,
            0.4,
            rec.get("title") or "",
            font_size=12,
            bold=True,
            color=DOWNY_TEAL,
        )

        _add_text_box(
            slide6,
            x + 0.15,
            y + 0.55,
            card_width - 0.3,
            0.85,
            rec.get("description") or "",
            font_size=9,
            color=WHITE,
        )

    # Footer
    _add_text_box(
        slide6,
        0.8,
        6.8,
        11,
        0.4,
        "Powered by Nova AI Suite | https://media-plan-generator.onrender.com",
        font_size=8,
        color=MUTED_TEXT,
        alignment=PP_ALIGN.CENTER,
    )

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ═══════════════════════════════════════════════════════════════════════════════
# 9. ORCHESTRATOR -- run_full_analysis
# ═══════════════════════════════════════════════════════════════════════════════


def run_full_analysis(
    company_name: str,
    competitors: List[str],
    industry: str = "general_entry_level",
    roles: Optional[List[str]] = None,
) -> Dict[str, Any]:
    """Single orchestrator function for the competitive intelligence dashboard.

    Thread-safe, never crashes (wraps all stages in try/except).
    Runs all analysis stages concurrently where possible.

    Args:
        company_name: Primary company to analyze
        competitors: List of competitor company names (1-5)
        industry: Canonical industry key
        roles: Optional list of target roles for benchmark context

    Returns:
        Full analysis dict containing all sections + generated brief.
    """
    start_time = time.time()

    result: Dict[str, Any] = {
        "status": "success",
        "company_name": company_name,
        "competitors": competitors,
        "industry": industry,
        "industry_label": INDUSTRY_LABEL_MAP.get(industry, industry),
        "roles": roles or [],
        "analysis_time_ms": 0,
        "errors": [],
    }

    # Validate inputs
    company_name = (company_name or "").strip()
    if not company_name:
        result["status"] = "error"
        result["errors"].append("Company name is required")
        return result

    competitors = [c.strip() for c in (competitors or []) if c and c.strip()]
    if not competitors:
        result["status"] = "error"
        result["errors"].append("At least one competitor is required")
        return result

    if len(competitors) > 5:
        competitors = competitors[:5]
        result["errors"].append("Truncated to 5 competitors (maximum)")

    # Normalize industry
    if industry not in INDUSTRY_LABEL_MAP:
        industry = "general_entry_level"
        result["industry"] = industry
        result["industry_label"] = INDUSTRY_LABEL_MAP[industry]

    # Run analysis stages concurrently
    with ThreadPoolExecutor(max_workers=4) as executor:
        f_competitors = executor.submit(
            _safe_call, analyze_competitors, company_name, competitors
        )
        f_hiring = executor.submit(
            _safe_call,
            compare_hiring_activity,
            company_name,
            competitors,
            industry,
            roles,
        )
        f_benchmarks = executor.submit(
            _safe_call, compare_ad_benchmarks, industry, roles
        )
        f_trends = executor.submit(
            _safe_call, get_market_trends, company_name, competitors
        )

        # Collect results
        try:
            result["competitor_analysis"] = f_competitors.result(timeout=45) or {}
        except Exception as exc:
            result["competitor_analysis"] = {}
            result["errors"].append(f"Competitor analysis failed: {exc}")

        try:
            result["hiring_activity"] = f_hiring.result(timeout=30) or {}
        except Exception as exc:
            result["hiring_activity"] = {}
            result["errors"].append(f"Hiring activity failed: {exc}")

        try:
            result["ad_benchmarks"] = f_benchmarks.result(timeout=15) or {}
        except Exception as exc:
            result["ad_benchmarks"] = {}
            result["errors"].append(f"Ad benchmarks failed: {exc}")

        try:
            result["market_trends"] = f_trends.result(timeout=20) or {}
        except Exception as exc:
            result["market_trends"] = {}
            result["errors"].append(f"Market trends failed: {exc}")

    # Generate brief with recommendations
    try:
        brief = generate_competitive_brief(result)
        result["brief"] = brief
    except Exception as exc:
        result["brief"] = {}
        result["errors"].append(f"Brief generation failed: {exc}")

    result["analysis_time_ms"] = int((time.time() - start_time) * 1000)

    if result["errors"]:
        result["status"] = "partial" if result.get("brief") else "error"

    return result


# ═══════════════════════════════════════════════════════════════════════════════
# COMPETITIVE THREAT ASSESSMENT (P1-3 / S49)
# ═══════════════════════════════════════════════════════════════════════════════


# Average monthly recruitment ad spend estimates by company size / tier.
# Source: Recruitment marketing competitive research 2025-2026.
_SPEND_TIERS: Dict[str, Dict[str, Any]] = {
    "enterprise": {
        "label": "Enterprise (10K+ employees)",
        "monthly_spend_low": 500_000,
        "monthly_spend_high": 5_000_000,
        "typical_channels": ["LinkedIn", "Indeed", "Programmatic", "Google Ads"],
    },
    "large": {
        "label": "Large (1K-10K employees)",
        "monthly_spend_low": 50_000,
        "monthly_spend_high": 500_000,
        "typical_channels": ["Indeed", "LinkedIn", "ZipRecruiter", "Programmatic"],
    },
    "mid": {
        "label": "Mid-Market (100-1K employees)",
        "monthly_spend_low": 10_000,
        "monthly_spend_high": 100_000,
        "typical_channels": ["Indeed", "ZipRecruiter", "Google Ads", "Craigslist"],
    },
    "small": {
        "label": "Small (<100 employees)",
        "monthly_spend_low": 1_000,
        "monthly_spend_high": 15_000,
        "typical_channels": ["Indeed", "Craigslist", "Facebook Jobs", "Google Ads"],
    },
}

# Industry-level competitive intensity heuristics (1-10 scale).
_INDUSTRY_THREAT_SCORES: Dict[str, float] = {
    "healthcare_medical": 8.5,
    "tech_engineering": 8.0,
    "logistics_supply_chain": 7.5,
    "blue_collar_trades": 7.0,
    "retail_consumer": 6.5,
    "finance_banking": 7.0,
    "hospitality_travel": 6.0,
    "construction_real_estate": 6.5,
    "education": 5.0,
    "general_entry_level": 5.5,
    "pharma_biotech": 7.5,
    "automotive": 6.5,
    "energy_utilities": 6.0,
    "manufacturing": 6.5,
    "food_beverage": 6.0,
    "insurance": 5.5,
    "telecommunications": 6.0,
    "aerospace_defense": 7.0,
    "media_entertainment": 5.5,
    "mental_health": 7.0,
    "legal_services": 5.5,
    "military_recruitment": 4.5,
    "maritime_marine": 5.0,
}

# Well-known top competitors by industry vertical.
_INDUSTRY_COMPETITORS: Dict[str, List[Dict[str, str]]] = {
    "healthcare_medical": [
        {
            "name": "HCA Healthcare",
            "size": "enterprise",
            "focus": "Hospitals, acute care",
        },
        {
            "name": "UnitedHealth Group",
            "size": "enterprise",
            "focus": "Insurance + care delivery",
        },
        {
            "name": "CVS Health",
            "size": "enterprise",
            "focus": "Pharmacy, clinics, insurance",
        },
        {
            "name": "Ascension Health",
            "size": "large",
            "focus": "Non-profit hospital system",
        },
    ],
    "tech_engineering": [
        {
            "name": "Amazon",
            "size": "enterprise",
            "focus": "Cloud, logistics, retail tech",
        },
        {"name": "Google", "size": "enterprise", "focus": "Search, cloud, AI"},
        {"name": "Meta", "size": "enterprise", "focus": "Social media, VR/AR"},
        {
            "name": "Microsoft",
            "size": "enterprise",
            "focus": "Cloud, enterprise software",
        },
    ],
    "logistics_supply_chain": [
        {
            "name": "Amazon Logistics",
            "size": "enterprise",
            "focus": "Last-mile delivery",
        },
        {"name": "FedEx", "size": "enterprise", "focus": "Express, ground, freight"},
        {
            "name": "UPS",
            "size": "enterprise",
            "focus": "Package delivery, supply chain",
        },
        {"name": "XPO Logistics", "size": "large", "focus": "Contract logistics, LTL"},
    ],
    "retail_consumer": [
        {"name": "Walmart", "size": "enterprise", "focus": "Mass retail, grocery"},
        {"name": "Amazon", "size": "enterprise", "focus": "E-commerce, warehouse"},
        {"name": "Target", "size": "enterprise", "focus": "General merchandise"},
        {"name": "Costco", "size": "enterprise", "focus": "Wholesale, membership"},
    ],
    "blue_collar_trades": [
        {
            "name": "Waste Management",
            "size": "large",
            "focus": "Environmental services",
        },
        {
            "name": "CenterPoint Energy",
            "size": "large",
            "focus": "Utilities, field services",
        },
        {"name": "ABM Industries", "size": "large", "focus": "Facility services"},
        {"name": "Cintas", "size": "large", "focus": "Uniforms, facility services"},
    ],
    "finance_banking": [
        {
            "name": "JPMorgan Chase",
            "size": "enterprise",
            "focus": "Banking, asset management",
        },
        {
            "name": "Bank of America",
            "size": "enterprise",
            "focus": "Consumer, commercial banking",
        },
        {"name": "Wells Fargo", "size": "enterprise", "focus": "Diversified banking"},
        {"name": "Goldman Sachs", "size": "enterprise", "focus": "Investment banking"},
    ],
}


def assess_competitive_threats(
    industry: str,
    role: str = "",
    locations: Optional[List[str]] = None,
) -> Dict[str, Any]:
    """Assess competitive threats for a given industry, role, and location set.

    Returns a structured threat assessment including:
    - threat_score (1-10)
    - top_competitors with spend estimates
    - market_position analysis
    - actionable recommendations

    Args:
        industry: Industry key (e.g., 'healthcare_medical', 'tech_engineering').
        role: Target role for context (e.g., 'Registered Nurse').
        locations: List of target locations (e.g., ['New York, NY', 'Dallas, TX']).

    Returns:
        Threat assessment dict with threat_score, competitors, position, recommendation.
    """
    locations = locations or []
    industry_lower = industry.lower().replace(" ", "_")

    # Normalize common industry aliases
    _aliases: Dict[str, str] = {
        "healthcare": "healthcare_medical",
        "medical": "healthcare_medical",
        "nursing": "healthcare_medical",
        "tech": "tech_engineering",
        "technology": "tech_engineering",
        "software": "tech_engineering",
        "it": "tech_engineering",
        "trucking": "logistics_supply_chain",
        "logistics": "logistics_supply_chain",
        "transportation": "logistics_supply_chain",
        "retail": "retail_consumer",
        "consumer": "retail_consumer",
        "warehouse": "blue_collar_trades",
        "manufacturing": "blue_collar_trades",
        "trades": "blue_collar_trades",
        "finance": "finance_banking",
        "banking": "finance_banking",
        "construction": "construction_real_estate",
    }
    industry_key = _aliases.get(industry_lower, industry_lower)

    # Base threat score from industry
    base_score = _INDUSTRY_THREAT_SCORES.get(industry_key, 5.5)

    # Location adjustment: major metros increase competition
    _high_competition_metros = {
        "new york",
        "los angeles",
        "chicago",
        "houston",
        "dallas",
        "san francisco",
        "seattle",
        "boston",
        "atlanta",
        "denver",
        "austin",
        "miami",
        "washington",
        "dc",
        "phoenix",
    }
    metro_boost = 0.0
    for loc in locations:
        loc_lower = loc.lower()
        if any(metro in loc_lower for metro in _high_competition_metros):
            metro_boost = max(metro_boost, 0.5)

    # Role adjustment: high-demand roles increase threat
    _high_demand_roles = [
        "nurse",
        "rn",
        "lpn",
        "cna",
        "developer",
        "engineer",
        "driver",
        "cdl",
        "mechanic",
        "electrician",
        "plumber",
        "data scientist",
        "warehouse",
        "picker",
        "packer",
        "technician",
    ]
    role_boost = 0.0
    role_lower = role.lower()
    if any(r in role_lower for r in _high_demand_roles):
        role_boost = 0.5

    threat_score = min(10.0, round(base_score + metro_boost + role_boost, 1))

    # Get competitors for industry
    competitors_raw = _INDUSTRY_COMPETITORS.get(industry_key, [])
    if not competitors_raw:
        # Fallback: use general competitors
        competitors_raw = [
            {
                "name": "Indeed Aggregators",
                "size": "enterprise",
                "focus": "Job board aggregation",
            },
            {
                "name": "LinkedIn Recruiter",
                "size": "enterprise",
                "focus": "Professional network",
            },
            {
                "name": "Staffing Agencies",
                "size": "large",
                "focus": "Temp-to-perm staffing",
            },
        ]

    # Enrich competitors with spend estimates
    top_competitors: List[Dict[str, Any]] = []
    for comp in competitors_raw[:5]:
        size_tier = comp.get("size", "mid")
        spend_info = _SPEND_TIERS.get(size_tier, _SPEND_TIERS["mid"])
        top_competitors.append(
            {
                "name": comp["name"],
                "size_tier": spend_info["label"],
                "focus": comp.get("focus", ""),
                "spend_estimate": f"${spend_info['monthly_spend_low']:,}-${spend_info['monthly_spend_high']:,}/mo",
                "typical_channels": spend_info["typical_channels"],
            }
        )

    # Market position assessment
    if threat_score >= 8.0:
        market_position = "highly_competitive"
        position_label = "Highly Competitive"
        position_detail = (
            f"The {INDUSTRY_LABEL_MAP.get(industry_key, industry)} market is "
            f"intensely competitive for talent. Expect aggressive bidding on "
            f"job boards and high CPC/CPA across channels."
        )
    elif threat_score >= 6.0:
        market_position = "moderately_competitive"
        position_label = "Moderately Competitive"
        position_detail = (
            f"The {INDUSTRY_LABEL_MAP.get(industry_key, industry)} market has "
            f"moderate competition. Strategic channel selection and compelling "
            f"creative can provide an edge."
        )
    else:
        market_position = "low_competition"
        position_label = "Lower Competition"
        position_detail = (
            f"The {INDUSTRY_LABEL_MAP.get(industry_key, industry)} market has "
            f"relatively lower hiring competition. Focus on efficiency and "
            f"cost optimization to maximize ROI."
        )

    # Generate actionable recommendations
    recommendations: List[str] = []

    if threat_score >= 7.5:
        recommendations.append(
            "Differentiate with salary transparency and benefits-first job postings "
            "(3.8x more applications when salary leads the posting)."
        )
        recommendations.append(
            "Diversify beyond top-2 job boards. Niche boards and programmatic "
            "channels often have 30-50% lower CPA in competitive markets."
        )
    if threat_score >= 6.0:
        recommendations.append(
            "Invest in employer branding content. Companies with strong employer "
            "brands see 50% more qualified applicants and 28% lower turnover."
        )
    if metro_boost > 0:
        recommendations.append(
            f"Major metro location(s) detected ({', '.join(locations[:3])}). "
            f"Consider geo-targeted campaigns with neighborhood-level targeting "
            f"and commute-time messaging."
        )
    if role_boost > 0:
        recommendations.append(
            f"High-demand role '{role}' detected. Speed-to-apply is critical -- "
            f"aim for <5 minute apply process (12.5% completion vs 3.5% for longer)."
        )
    recommendations.append(
        "Monitor competitor posting volume weekly using Nova's job volume tracker "
        "to identify windows of lower competition."
    )

    # Try enriching with live benchmark data
    cpc_context = {}
    if _HAS_BENCHMARK_REGISTRY:
        try:
            for ch in ["indeed", "linkedin", "google_search"]:
                bench = get_channel_benchmark(ch, industry=industry_key)
                if bench:
                    cpc_context[ch] = {
                        "cpc": bench.get("cpc"),
                        "cpa": bench.get("cpa"),
                    }
        except Exception as exc:
            logger.debug("Benchmark lookup for threat assessment failed: %s", exc)

    return {
        "threat_score": threat_score,
        "threat_level": position_label,
        "market_position": market_position,
        "industry": INDUSTRY_LABEL_MAP.get(industry_key, industry),
        "role": role or "General",
        "locations": locations,
        "top_competitors": top_competitors,
        "position_detail": position_detail,
        "recommendations": recommendations,
        "channel_benchmarks": cpc_context if cpc_context else None,
        "source": "Nova AI Competitive Intelligence Engine",
    }


# ═══════════════════════════════════════════════════════════════════════════════
# UTILITY HELPERS
# ═══════════════════════════════════════════════════════════════════════════════


def _safe_call(fn, *args, **kwargs):
    """Call a function, returning None on any exception."""
    try:
        return fn(*args, **kwargs)
    except Exception as exc:
        logger.warning("_safe_call(%s) failed: %s", fn.__name__, exc)
        return None


def _format_employee_count(count) -> str:
    """Format employee count for display."""
    if count is None:
        return "N/A"
    if isinstance(count, str):
        return count
    try:
        n = int(count)
    except (ValueError, TypeError):
        return str(count)
    if n >= 1_000_000:
        return f"{n / 1_000_000:.1f}M"
    if n >= 1_000:
        return f"{n / 1_000:.0f}K+"
    return str(n)


def get_industry_options() -> List[Dict[str, str]]:
    """Return industry options for the frontend dropdown."""
    return [{"value": key, "label": label} for key, label in INDUSTRY_LABEL_MAP.items()]
