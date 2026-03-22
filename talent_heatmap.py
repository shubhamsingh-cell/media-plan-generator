#!/usr/bin/env python3
"""
talent_heatmap.py -- Talent Supply Heat Map Backend

Uses the 100+ metros, 50 US states, and 40+ countries data from research.py
to create talent supply heat maps with:
  - Talent density scoring per location
  - Salary benchmarking across locations
  - Hiring competition / difficulty index
  - Cost-of-living adjusted comparisons
  - Optimal location recommendations
  - Side-by-side location comparison
  - Excel & PowerPoint report generation

Thread-safe, graceful degradation when APIs are unavailable.
All external data lookups use ThreadPoolExecutor for concurrent fetching.

Depends on (lazy-imported):
  - research (METRO_DATA, STATE_DATA, COUNTRY_DATA, detect_country, etc.)
  - shared_utils (INDUSTRY_LABEL_MAP)
  - api_enrichment (salary/demand APIs)
  - data_orchestrator (enrichment pipeline)
"""

from __future__ import annotations

import io
import logging
import math
import re
import time
import traceback
from collections import defaultdict
from concurrent.futures import ThreadPoolExecutor, as_completed
from datetime import datetime
from typing import Any, Dict, List, Optional, Tuple

logger = logging.getLogger(__name__)

# ═══════════════════════════════════════════════════════════════════════════════
# Lazy Imports -- graceful fallback when modules are unavailable
# ═══════════════════════════════════════════════════════════════════════════════

_research = None
_api_enrichment = None
_data_orchestrator = None
_trend_engine = None

_HAS_RESEARCH = False
_HAS_API = False
_HAS_ORCHESTRATOR = False
_HAS_TRENDS = False


def _lazy_research():
    global _research, _HAS_RESEARCH
    if _research is not None:
        return _research
    try:
        import research as _mod
        _research = _mod
        _HAS_RESEARCH = True
        return _mod
    except ImportError:
        logger.warning("research not available; heatmap will use fallback data")
        _HAS_RESEARCH = False
        return None


def _lazy_api():
    global _api_enrichment, _HAS_API
    if _api_enrichment is not None:
        return _api_enrichment
    try:
        import api_enrichment as _mod
        _api_enrichment = _mod
        _HAS_API = True
        return _mod
    except ImportError:
        logger.warning("api_enrichment not available; salary lookups will use fallbacks")
        _HAS_API = False
        return None


def _lazy_orchestrator():
    global _data_orchestrator, _HAS_ORCHESTRATOR
    if _data_orchestrator is not None:
        return _data_orchestrator
    try:
        import data_orchestrator as _mod
        _data_orchestrator = _mod
        _HAS_ORCHESTRATOR = True
        return _mod
    except ImportError:
        logger.warning("data_orchestrator not available; enrichment limited")
        _HAS_ORCHESTRATOR = False
        return None


def _lazy_trends():
    global _trend_engine, _HAS_TRENDS
    if _trend_engine is not None:
        return _trend_engine
    try:
        import trend_engine as _mod
        _trend_engine = _mod
        _HAS_TRENDS = True
        return _mod
    except ImportError:
        logger.warning("trend_engine not available; CPA benchmarks will use fallbacks")
        _HAS_TRENDS = False
        return None


# ═══════════════════════════════════════════════════════════════════════════════
# Industry label map (import from shared_utils or define fallback)
# ═══════════════════════════════════════════════════════════════════════════════

try:
    from shared_utils import INDUSTRY_LABEL_MAP
except ImportError:
    INDUSTRY_LABEL_MAP = {
        "healthcare_medical": "Healthcare & Medical",
        "blue_collar_trades": "Blue Collar / Skilled Trades",
        "tech_engineering": "Technology & Engineering",
        "general_entry_level": "General / Entry-Level",
        "finance_banking": "Finance & Banking",
        "retail_consumer": "Retail & Consumer",
        "logistics_supply_chain": "Logistics & Supply Chain",
        "hospitality_travel": "Hospitality & Travel",
        "construction_real_estate": "Construction & Real Estate",
        "education": "Education",
        "aerospace_defense": "Aerospace & Defense",
        "pharma_biotech": "Pharma & Biotech",
        "energy_utilities": "Energy & Utilities",
        "insurance": "Insurance",
        "telecommunications": "Telecommunications",
        "automotive": "Automotive & Manufacturing",
        "food_beverage": "Food & Beverage",
        "media_entertainment": "Media & Entertainment",
        "legal_services": "Legal Services",
        "mental_health": "Mental Health & Behavioral",
        "maritime_marine": "Maritime & Marine",
        "military_recruitment": "Military Recruitment",
    }

CANONICAL_INDUSTRIES = list(INDUSTRY_LABEL_MAP.keys())

# Max concurrent workers
_MAX_WORKERS = 10


# ═══════════════════════════════════════════════════════════════════════════════
# ROLE-BASED SALARY MULTIPLIERS
# ═══════════════════════════════════════════════════════════════════════════════

# Multiplier relative to metro median salary for common role families
ROLE_SALARY_MULTIPLIERS = {
    "software engineer": 1.35,
    "senior software engineer": 1.60,
    "staff software engineer": 1.85,
    "data scientist": 1.40,
    "data engineer": 1.35,
    "product manager": 1.45,
    "engineering manager": 1.65,
    "devops engineer": 1.30,
    "cloud engineer": 1.30,
    "machine learning engineer": 1.50,
    "ai engineer": 1.50,
    "frontend engineer": 1.25,
    "backend engineer": 1.30,
    "full stack developer": 1.25,
    "qa engineer": 1.10,
    "ux designer": 1.15,
    "ui designer": 1.10,
    "registered nurse": 0.95,
    "nurse practitioner": 1.25,
    "physician": 2.50,
    "medical assistant": 0.60,
    "pharmacist": 1.40,
    "physical therapist": 1.00,
    "dental hygienist": 0.85,
    "cna": 0.50,
    "lpn": 0.65,
    "occupational therapist": 1.05,
    "accountant": 0.95,
    "financial analyst": 1.10,
    "investment banker": 1.80,
    "actuary": 1.50,
    "compliance officer": 1.10,
    "underwriter": 1.00,
    "marketing manager": 1.20,
    "sales manager": 1.25,
    "account executive": 1.00,
    "business analyst": 1.05,
    "project manager": 1.10,
    "hr manager": 1.05,
    "recruiter": 0.85,
    "warehouse worker": 0.55,
    "forklift operator": 0.55,
    "truck driver": 0.70,
    "cdl driver": 0.75,
    "electrician": 0.80,
    "plumber": 0.78,
    "hvac technician": 0.75,
    "welder": 0.70,
    "machinist": 0.72,
    "carpenter": 0.68,
    "construction worker": 0.60,
    "mechanic": 0.65,
    "teacher": 0.70,
    "professor": 1.20,
    "paralegal": 0.75,
    "lawyer": 1.70,
    "executive assistant": 0.70,
    "administrative assistant": 0.55,
    "customer service representative": 0.50,
    "retail associate": 0.45,
    "cashier": 0.40,
    "restaurant manager": 0.70,
    "chef": 0.65,
    "line cook": 0.45,
    "security guard": 0.50,
}

# Industry-specific talent density multipliers (relative to baseline)
INDUSTRY_TALENT_DENSITY = {
    "tech_engineering": {"san_francisco": 2.5, "san_jose": 2.8, "seattle": 2.2, "austin": 2.0,
                         "new_york": 1.8, "boston": 1.9, "denver": 1.6, "raleigh": 1.7,
                         "los_angeles": 1.5, "chicago": 1.3, "atlanta": 1.4, "dallas": 1.3,
                         "portland": 1.5, "salt_lake_city": 1.4, "minneapolis": 1.2,
                         "washington_dc": 1.5, "phoenix": 1.2, "charlotte": 1.1},
    "healthcare_medical": {"boston": 2.2, "houston": 2.0, "nashville": 2.3, "new_york": 1.8,
                           "philadelphia": 1.9, "cleveland": 1.8, "rochester_mn": 2.5,
                           "baltimore": 1.7, "pittsburgh": 1.6, "minneapolis": 1.7,
                           "chicago": 1.5, "los_angeles": 1.4, "san_francisco": 1.3,
                           "dallas": 1.4, "atlanta": 1.5, "detroit": 1.3},
    "finance_banking": {"new_york": 3.0, "charlotte": 2.2, "chicago": 1.8, "boston": 1.7,
                        "san_francisco": 1.6, "washington_dc": 1.4, "dallas": 1.3,
                        "philadelphia": 1.3, "hartford": 1.5, "des_moines": 1.4,
                        "atlanta": 1.2, "minneapolis": 1.3, "omaha": 1.2},
    "logistics_supply_chain": {"memphis": 2.5, "louisville": 2.3, "dallas": 1.8,
                               "chicago": 1.7, "atlanta": 1.8, "indianapolis": 1.6,
                               "columbus": 1.5, "kansas_city": 1.5, "houston": 1.4,
                               "las_vegas": 1.2, "phoenix": 1.3, "jacksonville": 1.3},
    "aerospace_defense": {"washington_dc": 2.5, "huntsville": 2.3, "san_diego": 2.0,
                          "seattle": 1.8, "dallas": 1.6, "denver": 1.5, "los_angeles": 1.7,
                          "tucson": 1.4, "phoenix": 1.3, "oklahoma_city": 1.4,
                          "norfolk": 1.3, "baltimore": 1.2},
    "retail_consumer": {"new_york": 1.5, "los_angeles": 1.4, "chicago": 1.3, "dallas": 1.3,
                        "atlanta": 1.3, "houston": 1.2, "miami": 1.3, "phoenix": 1.2,
                        "las_vegas": 1.1, "tampa": 1.1, "orlando": 1.2},
    "energy_utilities": {"houston": 2.8, "dallas": 1.5, "denver": 1.6, "oklahoma_city": 1.8,
                         "pittsburgh": 1.3, "new_orleans": 1.4, "anchorage": 1.5,
                         "san_antonio": 1.2, "salt_lake_city": 1.2},
    "construction_real_estate": {"houston": 1.8, "dallas": 1.7, "phoenix": 1.8,
                                 "austin": 1.6, "las_vegas": 1.5, "denver": 1.5,
                                 "atlanta": 1.4, "tampa": 1.4, "nashville": 1.3,
                                 "charlotte": 1.3, "san_antonio": 1.2},
}

# Competition multipliers per industry
INDUSTRY_COMPETITION_FACTOR = {
    "tech_engineering": 1.4,
    "healthcare_medical": 1.3,
    "finance_banking": 1.2,
    "aerospace_defense": 1.3,
    "pharma_biotech": 1.35,
    "legal_services": 1.25,
    "logistics_supply_chain": 0.9,
    "blue_collar_trades": 0.85,
    "general_entry_level": 0.7,
    "retail_consumer": 0.75,
    "hospitality_travel": 0.8,
    "food_beverage": 0.7,
    "construction_real_estate": 0.9,
    "education": 0.85,
    "energy_utilities": 1.1,
    "automotive": 1.0,
    "insurance": 1.05,
    "telecommunications": 1.1,
    "media_entertainment": 1.15,
    "mental_health": 1.2,
    "maritime_marine": 0.95,
    "military_recruitment": 0.8,
}


# ═══════════════════════════════════════════════════════════════════════════════
# 1. LOCATION RESOLUTION
# ═══════════════════════════════════════════════════════════════════════════════

def _resolve_location(location: str) -> Dict[str, Any]:
    """Resolve a location string to its data from research.py.

    Tries METRO_DATA, STATE_DATA, and COUNTRY_DATA in that order.
    Returns dict with: type, key, name, data (the raw record).
    """
    if not location:
        return {"type": "unknown", "key": "", "name": location, "data": {}}

    res = _lazy_research()
    if not res:
        return {"type": "unknown", "key": location, "name": location, "data": {}}

    loc_lower = location.strip().lower()
    loc_clean = re.sub(r'[^a-z0-9\s]', '', loc_lower).strip()
    loc_underscore = re.sub(r'\s+', '_', loc_clean)

    # 1. Check METRO_DATA
    metro_data = getattr(res, "METRO_DATA", {})
    # Direct key match
    if loc_underscore in metro_data:
        d = metro_data[loc_underscore]
        return {"type": "metro", "key": loc_underscore,
                "name": d.get("metro_name", location.title()), "data": d}
    # Try without underscores and partial match
    for key, d in metro_data.items():
        metro_name_lower = d.get("metro_name", "").lower()
        key_clean = key.replace("_", " ")
        if loc_clean == key_clean or loc_clean in metro_name_lower:
            return {"type": "metro", "key": key,
                    "name": d.get("metro_name", key.replace("_", " ").title()), "data": d}

    # 2. Check COUNTRY_DATA (international locations)
    country_data = getattr(res, "COUNTRY_DATA", {})
    detect_country_fn = getattr(res, "detect_country", None)
    if detect_country_fn:
        country = detect_country_fn(location)
        if country and country in country_data:
            return {"type": "country", "key": country, "name": country,
                    "data": country_data[country]}

    # 3. Check STATE_DATA
    state_data = getattr(res, "STATE_DATA", {})
    # Check state abbreviation
    loc_upper = location.strip().upper()
    if loc_upper in state_data:
        d = state_data[loc_upper]
        return {"type": "state", "key": loc_upper,
                "name": d.get("name", loc_upper), "data": d}
    # Check state name
    for abbr, d in state_data.items():
        if d.get("name", "").lower() == loc_lower:
            return {"type": "state", "key": abbr, "name": d["name"], "data": d}

    # 4. Check metro data with city, state format ("Dallas, TX")
    parts = [p.strip() for p in location.split(",")]
    if len(parts) >= 2:
        city_part = re.sub(r'[^a-z0-9\s]', '', parts[0].lower()).strip()
        city_underscore = re.sub(r'\s+', '_', city_part)
        if city_underscore in metro_data:
            d = metro_data[city_underscore]
            return {"type": "metro", "key": city_underscore,
                    "name": d.get("metro_name", location.title()), "data": d}

    return {"type": "unknown", "key": location, "name": location, "data": {}}


def _get_all_metro_keys() -> List[str]:
    """Get all metro area keys from research.py."""
    res = _lazy_research()
    if res and hasattr(res, "METRO_DATA"):
        return list(res.METRO_DATA.keys())
    return []


# ═══════════════════════════════════════════════════════════════════════════════
# 2. TALENT DENSITY
# ═══════════════════════════════════════════════════════════════════════════════

def get_talent_density(role: str, locations: List[str],
                       industry: str = "general_entry_level") -> List[Dict[str, Any]]:
    """Calculate talent concentration per location for a given role.

    Uses population data, industry density multipliers, and metro employment
    composition to estimate relative talent availability.

    Returns list of dicts with: location, talent_density_score (0-100),
    talent_pool_estimate, population, top_industries, rating.
    """
    results = []
    role_lower = role.lower().strip() if role else ""

    for loc in locations:
        resolved = _resolve_location(loc)
        data = resolved["data"]
        loc_type = resolved["type"]
        loc_name = resolved["name"]
        loc_key = resolved["key"]

        if not data:
            results.append({
                "location": loc,
                "display_name": loc,
                "talent_density_score": 30,
                "talent_pool_estimate": "Unknown",
                "population": "N/A",
                "top_industries": "N/A",
                "rating": "Unknown",
            })
            continue

        # Base density from population
        pop_str = data.get("population", "0")
        pop_num = _parse_population(pop_str)
        base_density = min(100, max(10, math.log10(max(pop_num, 1000)) * 12))

        # Industry-specific density multiplier
        ind_density_map = INDUSTRY_TALENT_DENSITY.get(industry, {})
        ind_mult = ind_density_map.get(loc_key, 1.0)

        # Role relevance boost -- check if metro's top industries align with role
        top_ind = data.get("top_industries", "")
        role_boost = 1.0
        if _role_matches_industries(role_lower, top_ind):
            role_boost = 1.2

        # Major employers boost
        employers = data.get("major_employers", "")
        if employers and len(employers.split(",")) >= 4:
            role_boost *= 1.05

        # Unemployment factor (lower unemployment = harder to hire but indicates strong economy)
        unemp_str = data.get("unemployment", "4.0%")
        unemp = float(unemp_str.replace("%", "").strip()) if unemp_str else 4.0
        unemp_factor = 1.0 + (4.0 - unemp) * 0.02  # slight bonus for low unemployment

        density_score = min(100, max(5, base_density * ind_mult * role_boost * unemp_factor))

        # Estimate talent pool
        talent_pool = _estimate_talent_pool(pop_num, density_score, industry, role_lower)

        # Rating
        rating = _density_rating(density_score)

        results.append({
            "location": loc,
            "display_name": loc_name,
            "location_key": loc_key,
            "location_type": loc_type,
            "talent_density_score": round(density_score, 1),
            "talent_pool_estimate": talent_pool,
            "population": pop_str,
            "top_industries": top_ind,
            "major_employers": data.get("major_employers", "N/A"),
            "rating": rating,
        })

    # Normalize scores relative to the set
    if results:
        max_score = max(r["talent_density_score"] for r in results)
        if max_score > 0:
            for r in results:
                r["normalized_score"] = round(r["talent_density_score"] / max_score * 100, 1)

    return sorted(results, key=lambda x: x["talent_density_score"], reverse=True)


def _parse_population(pop_str: str) -> int:
    """Parse population string like '2.5M metro' to integer."""
    if not pop_str:
        return 0
    pop_clean = re.sub(r'[^0-9.MKBmkb]', '', str(pop_str))
    try:
        if 'B' in pop_clean.upper():
            return int(float(pop_clean.upper().replace('B', '')) * 1_000_000_000)
        elif 'M' in pop_clean.upper():
            return int(float(pop_clean.upper().replace('M', '')) * 1_000_000)
        elif 'K' in pop_clean.upper():
            return int(float(pop_clean.upper().replace('K', '')) * 1_000)
        return int(float(pop_clean))
    except (ValueError, TypeError):
        return 0


def _role_matches_industries(role: str, industries: str) -> bool:
    """Check if a role aligns with a location's top industries."""
    ind_lower = industries.lower()
    role_industry_map = {
        "software": ["technology", "tech", "engineering"],
        "nurse": ["healthcare", "medical", "health"],
        "doctor": ["healthcare", "medical", "health"],
        "physician": ["healthcare", "medical", "health"],
        "engineer": ["technology", "engineering", "manufacturing", "aerospace"],
        "financial": ["financial", "finance", "banking", "insurance"],
        "accountant": ["financial", "finance", "banking"],
        "teacher": ["education"],
        "driver": ["logistics", "transportation", "supply chain"],
        "warehouse": ["logistics", "manufacturing", "supply chain"],
        "construction": ["construction", "real estate"],
        "retail": ["retail", "consumer"],
        "chef": ["hospitality", "food", "tourism"],
        "mechanic": ["automotive", "manufacturing"],
        "sales": ["retail", "consumer", "technology"],
    }
    for keyword, ind_keywords in role_industry_map.items():
        if keyword in role:
            for ik in ind_keywords:
                if ik in ind_lower:
                    return True
    return False


def _estimate_talent_pool(population: int, density_score: float,
                          industry: str, role: str) -> str:
    """Estimate the approximate talent pool size for a role in a location."""
    # Base percentage of population in workforce (~60%)
    workforce = population * 0.60
    # Industry share (~5-15% depending on industry)
    industry_share = {
        "tech_engineering": 0.08, "healthcare_medical": 0.12,
        "finance_banking": 0.07, "retail_consumer": 0.10,
        "logistics_supply_chain": 0.06, "hospitality_travel": 0.09,
        "construction_real_estate": 0.06, "education": 0.08,
        "general_entry_level": 0.15, "blue_collar_trades": 0.08,
    }
    share = industry_share.get(industry, 0.07)
    # Role specificity factor (~1-10% of industry workers)
    role_specificity = 0.05
    if "senior" in role or "staff" in role or "manager" in role:
        role_specificity = 0.02
    elif "entry" in role or "junior" in role or "assistant" in role:
        role_specificity = 0.10

    pool = workforce * share * role_specificity * (density_score / 50)
    pool = max(50, pool)

    if pool >= 100000:
        return f"{pool / 1000:.0f}K+"
    elif pool >= 10000:
        return f"{pool / 1000:.1f}K"
    elif pool >= 1000:
        return f"{pool / 1000:.1f}K"
    return f"{int(pool)}"


def _density_rating(score: float) -> str:
    """Convert density score to a human-readable rating."""
    if score >= 75:
        return "Very High"
    elif score >= 55:
        return "High"
    elif score >= 40:
        return "Moderate"
    elif score >= 25:
        return "Low"
    return "Very Low"


# ═══════════════════════════════════════════════════════════════════════════════
# 3. SALARY MAP
# ═══════════════════════════════════════════════════════════════════════════════

def get_salary_map(role: str, locations: List[str],
                   industry: str = "general_entry_level") -> List[Dict[str, Any]]:
    """Get salary data per location for a given role.

    Uses METRO_DATA median_salary with role-specific multipliers
    and cost-of-living adjustments.

    Returns list of dicts with: location, estimated_salary, median_salary,
    coli, adjusted_salary, salary_range.
    """
    results = []
    role_lower = role.lower().strip() if role else ""

    # Find best matching role multiplier
    role_mult = 1.0
    for role_key, mult in ROLE_SALARY_MULTIPLIERS.items():
        if role_key in role_lower or role_lower in role_key:
            role_mult = mult
            break

    for loc in locations:
        resolved = _resolve_location(loc)
        data = resolved["data"]
        loc_name = resolved["name"]

        if not data:
            results.append({
                "location": loc,
                "display_name": loc,
                "estimated_salary": 0,
                "median_salary": 0,
                "coli": 100,
                "adjusted_salary": 0,
                "salary_range_low": 0,
                "salary_range_high": 0,
                "currency": "USD",
            })
            continue

        median = data.get("median_salary", 50000)
        coli = data.get("coli", 100)

        # Estimated salary for this role
        estimated = round(median * role_mult)

        # CoL-adjusted effective salary (what the salary buys)
        adjusted = round(estimated * (100 / max(coli, 50)))

        # Salary range (roughly +/- 20%)
        range_low = round(estimated * 0.80)
        range_high = round(estimated * 1.25)

        # Determine currency
        currency = "USD"
        if resolved["type"] == "country":
            currency = data.get("currency", "USD")

        results.append({
            "location": loc,
            "display_name": loc_name,
            "estimated_salary": estimated,
            "median_salary": median,
            "coli": coli,
            "adjusted_salary": adjusted,
            "salary_range_low": range_low,
            "salary_range_high": range_high,
            "currency": currency,
        })

    return sorted(results, key=lambda x: x["estimated_salary"], reverse=True)


# ═══════════════════════════════════════════════════════════════════════════════
# 4. COMPETITION INDEX
# ═══════════════════════════════════════════════════════════════════════════════

def get_competition_index(role: str, locations: List[str],
                          industry: str = "general_entry_level") -> List[Dict[str, Any]]:
    """Calculate hiring competition index per location.

    Higher score = more competition for candidates.
    Uses COLI, unemployment rate, industry factors, and metro size.

    Returns list of dicts with: location, competition_score (0-100),
    rating, factors.
    """
    results = []
    role_lower = role.lower().strip() if role else ""
    ind_factor = INDUSTRY_COMPETITION_FACTOR.get(industry, 1.0)

    for loc in locations:
        resolved = _resolve_location(loc)
        data = resolved["data"]
        loc_name = resolved["name"]
        loc_key = resolved["key"]

        if not data:
            results.append({
                "location": loc,
                "display_name": loc,
                "competition_score": 50,
                "rating": "Moderate",
                "factors": [],
            })
            continue

        factors = []

        # Factor 1: Cost of living (higher CoL = more competition for talent)
        coli = data.get("coli", 100)
        coli_score = min(30, max(0, (coli - 70) * 0.4))
        if coli >= 130:
            factors.append("High cost of living drives salary expectations up")
        elif coli <= 90:
            factors.append("Lower cost of living makes offers more attractive")

        # Factor 2: Unemployment (lower = harder to hire)
        unemp_str = data.get("unemployment", "4.0%")
        unemp = float(unemp_str.replace("%", "").strip()) if unemp_str else 4.0
        unemp_score = max(0, min(25, (5.0 - unemp) * 7))
        if unemp < 3.0:
            factors.append("Very low unemployment -- tight labor market")
        elif unemp > 5.0:
            factors.append("Higher unemployment -- larger available workforce")

        # Factor 3: Metro size (larger = more employers competing)
        pop = _parse_population(data.get("population", "0"))
        size_score = min(20, max(0, math.log10(max(pop, 1000)) * 3))
        if pop > 5_000_000:
            factors.append("Large metro with many competing employers")

        # Factor 4: Industry competition factor
        ind_score = ind_factor * 15
        if ind_factor >= 1.3:
            factors.append(f"High-competition industry ({INDUSTRY_LABEL_MAP.get(industry, industry)})")
        elif ind_factor <= 0.8:
            factors.append(f"Lower competition in {INDUSTRY_LABEL_MAP.get(industry, industry)}")

        # Factor 5: Industry talent density (high density in this location = more competition)
        ind_density = INDUSTRY_TALENT_DENSITY.get(industry, {}).get(loc_key, 1.0)
        density_score = min(10, (ind_density - 1.0) * 8)
        if ind_density >= 2.0:
            factors.append("Major talent hub -- many employers competing for same talent pool")

        total_score = min(100, max(5, coli_score + unemp_score + size_score + ind_score + density_score))

        rating = _competition_rating(total_score)

        results.append({
            "location": loc,
            "display_name": loc_name,
            "competition_score": round(total_score, 1),
            "rating": rating,
            "factors": factors[:3],  # Top 3 factors
            "unemployment": unemp_str,
            "coli": coli,
        })

    return sorted(results, key=lambda x: x["competition_score"], reverse=True)


def _competition_rating(score: float) -> str:
    """Convert competition score to rating."""
    if score >= 75:
        return "Very High"
    elif score >= 55:
        return "High"
    elif score >= 40:
        return "Moderate"
    elif score >= 25:
        return "Low"
    return "Very Low"


# ═══════════════════════════════════════════════════════════════════════════════
# 5. COST OF LIVING
# ═══════════════════════════════════════════════════════════════════════════════

def get_cost_of_living_factor(locations: List[str]) -> List[Dict[str, Any]]:
    """Get cost-of-living index and factors per location.

    Returns list of dicts with: location, coli, rating, median_salary,
    housing_note.
    """
    results = []

    for loc in locations:
        resolved = _resolve_location(loc)
        data = resolved["data"]
        loc_name = resolved["name"]

        coli = data.get("coli", 100) if data else 100
        median = data.get("median_salary", data.get("median_household_income", 50000)) if data else 50000

        if coli >= 150:
            rating = "Very High"
            note = "Extremely expensive market -- expect premium salary demands"
        elif coli >= 120:
            rating = "High"
            note = "Above-average costs -- competitive offers needed"
        elif coli >= 100:
            rating = "Average"
            note = "National average cost of living"
        elif coli >= 85:
            rating = "Below Average"
            note = "Lower costs -- salary goes further"
        else:
            rating = "Low"
            note = "Very affordable market -- strong purchasing power"

        results.append({
            "location": loc,
            "display_name": loc_name,
            "coli": coli,
            "coli_rating": rating,
            "median_salary": median,
            "housing_note": note,
            "relative_to_national": f"{coli - 100:+d}%" if coli != 100 else "National average",
        })

    return sorted(results, key=lambda x: x["coli"], reverse=True)


# ═══════════════════════════════════════════════════════════════════════════════
# 6. HIRING DIFFICULTY (Composite Score)
# ═══════════════════════════════════════════════════════════════════════════════

def calculate_hiring_difficulty(role: str, location: str,
                                industry: str = "general_entry_level") -> Dict[str, Any]:
    """Calculate composite hiring difficulty score for a role at a location.

    Combines talent density, competition, salary demands, and CoL into
    a single 0-100 difficulty score. Higher = harder to hire.
    """
    density = get_talent_density(role, [location], industry)
    competition = get_competition_index(role, [location], industry)
    salary = get_salary_map(role, [location], industry)
    col = get_cost_of_living_factor([location])

    d = density[0] if density else {}
    c = competition[0] if competition else {}
    s = salary[0] if salary else {}
    l = col[0] if col else {}

    # Invert density (low density = harder to hire)
    density_difficulty = max(0, 100 - d.get("talent_density_score", 50))

    # Competition directly maps to difficulty
    competition_difficulty = c.get("competition_score", 50)

    # Salary premium (higher salary = harder/more expensive)
    estimated_salary = s.get("estimated_salary", 50000)
    salary_difficulty = min(100, max(0, (estimated_salary - 40000) / 1500))

    # CoL factor
    coli = l.get("coli", 100)
    col_difficulty = min(100, max(0, (coli - 70) * 1.2))

    # Weighted composite
    difficulty = (
        density_difficulty * 0.30 +
        competition_difficulty * 0.30 +
        salary_difficulty * 0.20 +
        col_difficulty * 0.20
    )
    difficulty = round(min(100, max(0, difficulty)), 1)

    if difficulty >= 75:
        rating = "Very Hard"
        color = "#DC2626"
        recommendation = "Consider remote hiring, relocation packages, or alternative locations with lower difficulty."
    elif difficulty >= 55:
        rating = "Hard"
        color = "#F97316"
        recommendation = "Competitive offers and strong employer branding are essential. Consider expanding search radius."
    elif difficulty >= 40:
        rating = "Moderate"
        color = "#EAB308"
        recommendation = "Standard recruitment approaches should work. Focus on quality job postings and multi-channel distribution."
    elif difficulty >= 25:
        rating = "Easy"
        color = "#22C55E"
        recommendation = "Good talent availability. Focus on efficient screening and fast hiring process."
    else:
        rating = "Very Easy"
        color = "#16A34A"
        recommendation = "Abundant talent supply. Optimize for quality and cultural fit."

    return {
        "location": location,
        "display_name": d.get("display_name", location),
        "role": role,
        "industry": industry,
        "difficulty_score": difficulty,
        "difficulty_rating": rating,
        "difficulty_color": color,
        "recommendation": recommendation,
        "components": {
            "talent_scarcity": round(density_difficulty, 1),
            "competition": round(competition_difficulty, 1),
            "salary_pressure": round(salary_difficulty, 1),
            "cost_of_living": round(col_difficulty, 1),
        },
        "details": {
            "talent_density_score": d.get("talent_density_score", 0),
            "talent_pool": d.get("talent_pool_estimate", "N/A"),
            "competition_score": c.get("competition_score", 0),
            "estimated_salary": estimated_salary,
            "coli": coli,
        },
    }


# ═══════════════════════════════════════════════════════════════════════════════
# 7. FIND OPTIMAL LOCATIONS
# ═══════════════════════════════════════════════════════════════════════════════

def find_optimal_locations(role: str, budget: float, num_hires: int,
                           industry: str = "general_entry_level",
                           max_results: int = 10) -> List[Dict[str, Any]]:
    """Find the best locations to hire for a given role and budget.

    Ranks all metro areas by a composite score combining:
    talent availability, hiring difficulty, cost efficiency, and budget fit.

    Returns top N locations with scores and recommendations.
    """
    all_metros = _get_all_metro_keys()
    if not all_metros:
        return []

    # Score all locations concurrently
    scored = []
    budget_per_hire = budget / max(num_hires, 1) if budget > 0 else 50000

    with ThreadPoolExecutor(max_workers=_MAX_WORKERS) as executor:
        futures = {
            executor.submit(_score_location, metro, role, industry, budget_per_hire): metro
            for metro in all_metros
        }
        for future in as_completed(futures):
            try:
                result = future.result(timeout=10)
                if result:
                    scored.append(result)
            except Exception:
                pass

    # Sort by composite score descending
    scored.sort(key=lambda x: x["composite_score"], reverse=True)

    # Take top N
    top = scored[:max_results]

    # Add ranking and reasoning
    for i, loc in enumerate(top):
        loc["rank"] = i + 1
        loc["reasoning"] = _generate_location_reasoning(loc, role, industry)

    return top


def _score_location(metro_key: str, role: str, industry: str,
                    budget_per_hire: float) -> Optional[Dict[str, Any]]:
    """Score a single metro area for optimal location ranking."""
    res = _lazy_research()
    if not res:
        return None

    metro_data = getattr(res, "METRO_DATA", {})
    data = metro_data.get(metro_key)
    if not data:
        return None

    metro_name = data.get("metro_name", metro_key.replace("_", " ").title())
    state = data.get("state", "")
    median_salary = data.get("median_salary", 50000)
    coli = data.get("coli", 100)
    pop = _parse_population(data.get("population", "0"))

    # Role-adjusted salary
    role_lower = role.lower().strip()
    role_mult = 1.0
    for rk, mult in ROLE_SALARY_MULTIPLIERS.items():
        if rk in role_lower or role_lower in rk:
            role_mult = mult
            break
    estimated_salary = round(median_salary * role_mult)

    # Budget fit (lower salary relative to budget = better fit)
    if budget_per_hire > 0:
        budget_fit = min(100, max(0, (1 - estimated_salary / (budget_per_hire * 3)) * 100))
    else:
        budget_fit = 50

    # Talent density
    ind_density = INDUSTRY_TALENT_DENSITY.get(industry, {}).get(metro_key, 1.0)
    density_score = min(100, max(10, math.log10(max(pop, 1000)) * 12 * ind_density))

    # Cost efficiency (lower CoL = more efficient)
    cost_efficiency = min(100, max(0, (200 - coli) * 0.6))

    # Competition (inverted -- lower competition is better)
    ind_factor = INDUSTRY_COMPETITION_FACTOR.get(industry, 1.0)
    unemp_str = data.get("unemployment", "4.0%")
    unemp = float(unemp_str.replace("%", "").strip()) if unemp_str else 4.0
    competition_inv = max(0, min(100, unemp * 15 + (100 - coli * 0.5)))

    # Composite score
    composite = (
        density_score * 0.30 +
        cost_efficiency * 0.25 +
        budget_fit * 0.25 +
        competition_inv * 0.20
    )

    return {
        "metro_key": metro_key,
        "location": metro_name,
        "state": state,
        "composite_score": round(composite, 1),
        "talent_density_score": round(density_score, 1),
        "cost_efficiency_score": round(cost_efficiency, 1),
        "budget_fit_score": round(budget_fit, 1),
        "competition_score": round(100 - competition_inv, 1),
        "estimated_salary": estimated_salary,
        "median_salary": median_salary,
        "coli": coli,
        "population": data.get("population", "N/A"),
        "unemployment": unemp_str,
        "major_employers": data.get("major_employers", ""),
        "top_industries": data.get("top_industries", ""),
    }


def _generate_location_reasoning(loc: Dict[str, Any], role: str, industry: str) -> str:
    """Generate a brief reasoning for why this location is recommended."""
    reasons = []
    if loc.get("talent_density_score", 0) >= 60:
        reasons.append("strong talent pool")
    if loc.get("cost_efficiency_score", 0) >= 60:
        reasons.append("cost-effective")
    if loc.get("budget_fit_score", 0) >= 60:
        reasons.append("good budget fit")
    if loc.get("competition_score", 0) <= 40:
        reasons.append("lower hiring competition")

    if not reasons:
        reasons.append("balanced across all factors")

    return f"Recommended for {role}: {', '.join(reasons)}."


# ═══════════════════════════════════════════════════════════════════════════════
# 8. COMPARE LOCATIONS
# ═══════════════════════════════════════════════════════════════════════════════

def compare_locations(role: str, locations: List[str],
                      industry: str = "general_entry_level") -> Dict[str, Any]:
    """Side-by-side comparison of multiple locations for a given role.

    Returns comprehensive comparison dict with all metrics aligned.
    """
    if not locations:
        return {"error": "At least one location is required", "locations": []}

    # Fetch all metrics concurrently
    with ThreadPoolExecutor(max_workers=4) as executor:
        f_density = executor.submit(get_talent_density, role, locations, industry)
        f_salary = executor.submit(get_salary_map, role, locations, industry)
        f_competition = executor.submit(get_competition_index, role, locations, industry)
        f_col = executor.submit(get_cost_of_living_factor, locations)

        density_data = _safe_result(f_density) or []
        salary_data = _safe_result(f_salary) or []
        competition_data = _safe_result(f_competition) or []
        col_data = _safe_result(f_col) or []

    # Index by location for easy lookup
    density_map = {d["location"]: d for d in density_data}
    salary_map = {s["location"]: s for s in salary_data}
    comp_map = {c["location"]: c for c in competition_data}
    col_map = {c["location"]: c for c in col_data}

    comparisons = []
    for loc in locations:
        d = density_map.get(loc, {})
        s = salary_map.get(loc, {})
        c = comp_map.get(loc, {})
        l = col_map.get(loc, {})

        # Calculate overall recommendation score
        density_score = d.get("talent_density_score", 50)
        comp_score = c.get("competition_score", 50)
        salary = s.get("estimated_salary", 50000)
        coli = l.get("coli", 100)

        # Higher density + lower competition + lower cost = better
        rec_score = (density_score * 0.4 +
                     (100 - comp_score) * 0.3 +
                     (200 - coli) * 0.15 +
                     max(0, 100 - salary / 1500) * 0.15)
        rec_score = round(min(100, max(0, rec_score)), 1)

        comparisons.append({
            "location": loc,
            "display_name": d.get("display_name", loc),
            "recommendation_score": rec_score,
            "talent_density": {
                "score": density_score,
                "pool": d.get("talent_pool_estimate", "N/A"),
                "rating": d.get("rating", "Unknown"),
            },
            "salary": {
                "estimated": s.get("estimated_salary", 0),
                "median": s.get("median_salary", 0),
                "range_low": s.get("salary_range_low", 0),
                "range_high": s.get("salary_range_high", 0),
                "adjusted": s.get("adjusted_salary", 0),
            },
            "competition": {
                "score": comp_score,
                "rating": c.get("rating", "Unknown"),
                "factors": c.get("factors", []),
            },
            "cost_of_living": {
                "index": coli,
                "rating": l.get("coli_rating", "Unknown"),
                "vs_national": l.get("relative_to_national", "N/A"),
            },
            "population": d.get("population", "N/A"),
            "top_industries": d.get("top_industries", "N/A"),
            "major_employers": d.get("major_employers", "N/A"),
        })

    # Sort by recommendation score
    comparisons.sort(key=lambda x: x["recommendation_score"], reverse=True)

    # Mark the best option
    if comparisons:
        comparisons[0]["is_best"] = True

    return {
        "role": role,
        "industry": industry,
        "industry_label": INDUSTRY_LABEL_MAP.get(industry, industry),
        "locations": comparisons,
        "best_location": comparisons[0]["location"] if comparisons else None,
    }


def _safe_result(future):
    """Get future result safely."""
    try:
        return future.result(timeout=15)
    except Exception as exc:
        logger.warning("Future failed: %s", exc)
        return None


# ═══════════════════════════════════════════════════════════════════════════════
# 9. EXCEL REPORT GENERATION
# ═══════════════════════════════════════════════════════════════════════════════

def generate_heatmap_excel(analysis: Dict[str, Any]) -> bytes:
    """Generate a Talent Supply Heatmap Excel workbook.

    Sheets:
      1. Location Rankings -- sorted table with all metrics
      2. Salary Comparison -- salary data across locations
      3. Hiring Difficulty -- composite difficulty scores
      4. Optimal Locations -- top recommended locations
      5. Cost Analysis -- total hiring cost estimates

    Uses Sapphire Blue palette (Calibri, column B start).
    """
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    # Design tokens -- Sapphire Blue palette
    NAVY = "0F172A"
    SAPPHIRE = "2563EB"
    BLUE_LIGHT = "DBEAFE"
    BLUE_PALE = "EFF6FF"
    GREEN = "16A34A"
    ORANGE = "F97316"
    RED = "DC2626"
    WARM_GRAY = "E7E5E4"
    MUTED = "78716C"

    hdr_fill = PatternFill(start_color=NAVY, end_color=NAVY, fill_type="solid")
    hdr_font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
    accent_fill = PatternFill(start_color=BLUE_PALE, end_color=BLUE_PALE, fill_type="solid")
    data_font = Font(name="Calibri", size=10, color="1E293B")
    bold_font = Font(name="Calibri", size=10, bold=True, color="1E293B")
    title_font = Font(name="Calibri", size=14, bold=True, color=NAVY)
    subtitle_font = Font(name="Calibri", size=11, color=MUTED)
    thin_border = Border(bottom=Side(style="thin", color=WARM_GRAY))
    wrap_align = Alignment(wrap_text=True, vertical="top")
    center_align = Alignment(horizontal="center", vertical="center")

    green_font = Font(name="Calibri", size=10, bold=True, color=GREEN)
    orange_font = Font(name="Calibri", size=10, bold=True, color=ORANGE)
    red_font = Font(name="Calibri", size=10, bold=True, color=RED)

    wb = Workbook()

    role = analysis.get("role", "")
    industry = analysis.get("industry", "")
    industry_label = analysis.get("industry_label", INDUSTRY_LABEL_MAP.get(industry, industry))

    # ── Sheet 1: Location Rankings ──
    ws1 = wb.active
    ws1.title = "Location Rankings"
    ws1.sheet_properties.tabColor = SAPPHIRE

    ws1.merge_cells("B2:I2")
    ws1["B2"] = f"Talent Supply Heatmap: {role}"
    ws1["B2"].font = title_font

    ws1.merge_cells("B3:I3")
    ws1["B3"] = f"Industry: {industry_label} | Generated {datetime.utcnow().strftime('%Y-%m-%d')} | Powered by Nova AI Suite"
    ws1["B3"].font = subtitle_font

    headers1 = ["Location", "Talent Density", "Talent Pool", "Competition",
                "Estimated Salary", "Cost of Living", "Hiring Difficulty", "Recommendation"]
    row = 5
    for i, h in enumerate(headers1):
        cell = ws1.cell(row=row, column=i + 2, value=h)
        cell.font = hdr_font
        cell.fill = hdr_fill
        cell.alignment = center_align

    comparison = analysis.get("comparison", {})
    locations = comparison.get("locations", [])
    row = 6
    for loc in locations:
        ws1.cell(row=row, column=2, value=loc.get("display_name", loc.get("location", ""))).font = bold_font
        td = loc.get("talent_density", {})
        ws1.cell(row=row, column=3, value=f"{td.get('score', 0)}/100 ({td.get('rating', '')})").font = data_font
        ws1.cell(row=row, column=4, value=td.get("pool", "N/A")).font = data_font
        comp = loc.get("competition", {})
        ws1.cell(row=row, column=5, value=f"{comp.get('score', 0)}/100 ({comp.get('rating', '')})").font = data_font
        sal = loc.get("salary", {})
        ws1.cell(row=row, column=6, value=f"${sal.get('estimated', 0):,}").font = data_font
        col = loc.get("cost_of_living", {})
        ws1.cell(row=row, column=7, value=f"{col.get('index', 100)} ({col.get('rating', '')})").font = data_font
        rec = loc.get("recommendation_score", 0)
        cell = ws1.cell(row=row, column=8, value=f"{rec}/100")
        cell.font = green_font if rec >= 60 else (orange_font if rec >= 40 else red_font)
        ws1.cell(row=row, column=9, value="Best" if loc.get("is_best") else "").font = green_font
        for c in range(2, 10):
            ws1.cell(row=row, column=c).border = thin_border
        row += 1

    ws1.column_dimensions["A"].width = 3
    ws1.column_dimensions["B"].width = 32
    for c in "CDEFGHI":
        ws1.column_dimensions[c].width = 20

    # ── Sheet 2: Salary Comparison ──
    ws2 = wb.create_sheet("Salary Comparison")
    ws2.sheet_properties.tabColor = SAPPHIRE

    ws2.merge_cells("B2:G2")
    ws2["B2"] = f"Salary Comparison: {role}"
    ws2["B2"].font = title_font

    headers2 = ["Location", "Estimated Salary", "Median Area Salary", "CoL Index",
                "CoL-Adjusted Salary", "Salary Range"]
    row = 4
    for i, h in enumerate(headers2):
        cell = ws2.cell(row=row, column=i + 2, value=h)
        cell.font = hdr_font
        cell.fill = hdr_fill
        cell.alignment = center_align

    salary_data = analysis.get("salary_map", [])
    row = 5
    for s in salary_data:
        ws2.cell(row=row, column=2, value=s.get("display_name", s.get("location", ""))).font = bold_font
        ws2.cell(row=row, column=3, value=f"${s.get('estimated_salary', 0):,}").font = data_font
        ws2.cell(row=row, column=4, value=f"${s.get('median_salary', 0):,}").font = data_font
        ws2.cell(row=row, column=5, value=s.get("coli", 100)).font = data_font
        ws2.cell(row=row, column=6, value=f"${s.get('adjusted_salary', 0):,}").font = data_font
        ws2.cell(row=row, column=7, value=f"${s.get('salary_range_low', 0):,} - ${s.get('salary_range_high', 0):,}").font = data_font
        for c in range(2, 8):
            ws2.cell(row=row, column=c).border = thin_border
        row += 1

    ws2.column_dimensions["A"].width = 3
    ws2.column_dimensions["B"].width = 32
    for c in "CDEFG":
        ws2.column_dimensions[c].width = 22

    # ── Sheet 3: Hiring Difficulty ──
    ws3 = wb.create_sheet("Hiring Difficulty")
    ws3.sheet_properties.tabColor = SAPPHIRE

    ws3.merge_cells("B2:H2")
    ws3["B2"] = f"Hiring Difficulty Analysis: {role}"
    ws3["B2"].font = title_font

    headers3 = ["Location", "Difficulty Score", "Rating", "Talent Scarcity",
                "Competition", "Salary Pressure", "Cost of Living"]
    row = 4
    for i, h in enumerate(headers3):
        cell = ws3.cell(row=row, column=i + 2, value=h)
        cell.font = hdr_font
        cell.fill = hdr_fill
        cell.alignment = center_align

    difficulties = analysis.get("difficulties", [])
    row = 5
    for d in difficulties:
        ws3.cell(row=row, column=2, value=d.get("display_name", d.get("location", ""))).font = bold_font
        score = d.get("difficulty_score", 0)
        cell = ws3.cell(row=row, column=3, value=f"{score}/100")
        cell.font = red_font if score >= 60 else (orange_font if score >= 40 else green_font)
        ws3.cell(row=row, column=4, value=d.get("difficulty_rating", "")).font = data_font
        comps = d.get("components", {})
        ws3.cell(row=row, column=5, value=f"{comps.get('talent_scarcity', 0)}/100").font = data_font
        ws3.cell(row=row, column=6, value=f"{comps.get('competition', 0)}/100").font = data_font
        ws3.cell(row=row, column=7, value=f"{comps.get('salary_pressure', 0)}/100").font = data_font
        ws3.cell(row=row, column=8, value=f"{comps.get('cost_of_living', 0)}/100").font = data_font
        for c in range(2, 9):
            ws3.cell(row=row, column=c).border = thin_border
        row += 1

    ws3.column_dimensions["A"].width = 3
    ws3.column_dimensions["B"].width = 32
    for c in "CDEFGH":
        ws3.column_dimensions[c].width = 18

    # ── Sheet 4: Optimal Locations ──
    ws4 = wb.create_sheet("Optimal Locations")
    ws4.sheet_properties.tabColor = SAPPHIRE

    ws4.merge_cells("B2:I2")
    ws4["B2"] = f"Optimal Hiring Locations: {role}"
    ws4["B2"].font = title_font

    headers4 = ["Rank", "Location", "Score", "Talent Density", "Cost Efficiency",
                "Budget Fit", "Est. Salary", "Reasoning"]
    row = 4
    for i, h in enumerate(headers4):
        cell = ws4.cell(row=row, column=i + 2, value=h)
        cell.font = hdr_font
        cell.fill = hdr_fill
        cell.alignment = center_align

    optimal = analysis.get("optimal_locations", [])
    row = 5
    for loc in optimal:
        ws4.cell(row=row, column=2, value=f"#{loc.get('rank', '')}").font = bold_font
        ws4.cell(row=row, column=3, value=loc.get("location", "")).font = bold_font
        cell = ws4.cell(row=row, column=4, value=f"{loc.get('composite_score', 0)}/100")
        cell.font = green_font
        ws4.cell(row=row, column=5, value=f"{loc.get('talent_density_score', 0)}/100").font = data_font
        ws4.cell(row=row, column=6, value=f"{loc.get('cost_efficiency_score', 0)}/100").font = data_font
        ws4.cell(row=row, column=7, value=f"{loc.get('budget_fit_score', 0)}/100").font = data_font
        ws4.cell(row=row, column=8, value=f"${loc.get('estimated_salary', 0):,}").font = data_font
        ws4.cell(row=row, column=9, value=loc.get("reasoning", "")).font = data_font
        ws4.cell(row=row, column=9).alignment = wrap_align
        for c in range(2, 10):
            ws4.cell(row=row, column=c).border = thin_border
        row += 1

    ws4.column_dimensions["A"].width = 3
    ws4.column_dimensions["B"].width = 8
    ws4.column_dimensions["C"].width = 30
    for c in "DEFGH":
        ws4.column_dimensions[c].width = 18
    ws4.column_dimensions["I"].width = 40

    # ── Sheet 5: Cost Analysis ──
    ws5 = wb.create_sheet("Cost Analysis")
    ws5.sheet_properties.tabColor = SAPPHIRE

    ws5.merge_cells("B2:G2")
    ws5["B2"] = f"Total Hiring Cost Estimate: {role}"
    ws5["B2"].font = title_font

    num_hires = analysis.get("num_hires", 5)
    ws5.merge_cells("B3:G3")
    ws5["B3"] = f"Based on {num_hires} hire(s) | Budget: ${analysis.get('budget', 0):,}"
    ws5["B3"].font = subtitle_font

    headers5 = ["Location", "Salary/Hire", "Recruiting Cost Est.", "Total Cost/Hire",
                "Total for All Hires", "Budget Fit"]
    row = 5
    for i, h in enumerate(headers5):
        cell = ws5.cell(row=row, column=i + 2, value=h)
        cell.font = hdr_font
        cell.fill = hdr_fill
        cell.alignment = center_align

    row = 6
    budget = analysis.get("budget", 0)
    for loc in locations:
        sal = loc.get("salary", {})
        estimated = sal.get("estimated", 50000)
        # Recruiting cost ~15-25% of salary
        recruiting_cost = round(estimated * 0.20)
        total_per_hire = estimated + recruiting_cost
        total_all = total_per_hire * num_hires

        ws5.cell(row=row, column=2, value=loc.get("display_name", "")).font = bold_font
        ws5.cell(row=row, column=3, value=f"${estimated:,}").font = data_font
        ws5.cell(row=row, column=4, value=f"${recruiting_cost:,}").font = data_font
        ws5.cell(row=row, column=5, value=f"${total_per_hire:,}").font = data_font
        ws5.cell(row=row, column=6, value=f"${total_all:,}").font = bold_font
        fit = "Within Budget" if total_all <= budget else "Over Budget"
        cell = ws5.cell(row=row, column=7, value=fit)
        cell.font = green_font if fit == "Within Budget" else red_font
        for c in range(2, 8):
            ws5.cell(row=row, column=c).border = thin_border
        row += 1

    ws5.column_dimensions["A"].width = 3
    ws5.column_dimensions["B"].width = 32
    for c in "CDEFG":
        ws5.column_dimensions[c].width = 22

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ═══════════════════════════════════════════════════════════════════════════════
# 10. POWERPOINT REPORT GENERATION
# ═══════════════════════════════════════════════════════════════════════════════

def generate_heatmap_ppt(analysis: Dict[str, Any]) -> bytes:
    """Generate a Talent Supply Heatmap PowerPoint presentation.

    Uses branding: Port Gore (#202058), Blue Violet (#5A54BD),
    Downy teal (#6BB3CD).
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        logger.error("python-pptx not available for PPT generation")
        return b""

    # Color constants -- Nova AI Suite branding
    PORT_GORE = RGBColor(0x20, 0x20, 0x58)
    BLUE_VIOLET = RGBColor(0x5A, 0x54, 0xBD)
    DOWNY_TEAL = RGBColor(0x6B, 0xB3, 0xCD)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    MUTED_TEXT = RGBColor(0x94, 0xA3, 0xB8)
    GREEN = RGBColor(0x22, 0xC5, 0x5E)
    ORANGE = RGBColor(0xF9, 0x73, 0x16)
    RED = RGBColor(0xDC, 0x26, 0x26)
    CARD_BG = RGBColor(0x1A, 0x1A, 0x30)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    role = analysis.get("role", "")
    industry_label = analysis.get("industry_label", "")

    def _add_bg(slide, color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_text_box(slide, left, top, width, height, text,
                      font_size=12, bold=False, color=WHITE,
                      alignment=PP_ALIGN.LEFT):
        from pptx.util import Inches, Pt
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(text)
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = alignment
        return txBox

    def _add_shape(slide, left, top, width, height, color):
        from pptx.util import Inches
        from pptx.enum.shapes import MSO_SHAPE
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    # ── Slide 1: Title ──
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide1, RGBColor(0x0F, 0x0F, 0x1E))

    _add_text_box(slide1, 1.5, 2.0, 10, 1.0,
                  "Talent Supply Heat Map", font_size=36, bold=True, color=DOWNY_TEAL)
    _add_text_box(slide1, 1.5, 3.2, 10, 0.6,
                  f"Role: {role} | Industry: {industry_label}", font_size=18, color=WHITE)
    _add_text_box(slide1, 1.5, 4.0, 10, 0.5,
                  f"Generated {datetime.utcnow().strftime('%B %d, %Y')} | Powered by Nova AI",
                  font_size=12, color=MUTED_TEXT)
    _add_text_box(slide1, 1.5, 6.5, 10, 0.4,
                  "Nova AI Suite | https://media-plan-generator.onrender.com",
                  font_size=10, color=MUTED_TEXT, alignment=PP_ALIGN.LEFT)

    # ── Slide 2: Location Rankings ──
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide2, RGBColor(0x0F, 0x0F, 0x1E))
    _add_text_box(slide2, 0.8, 0.4, 8, 0.6,
                  "Location Rankings", font_size=24, bold=True, color=WHITE)

    comparison = analysis.get("comparison", {})
    locations = comparison.get("locations", [])

    y = 1.3
    for i, loc in enumerate(locations[:8]):
        x = 0.8
        card_w = 11.5
        _add_shape(slide2, x, y, card_w, 0.65, CARD_BG)

        name = loc.get("display_name", loc.get("location", ""))
        td = loc.get("talent_density", {})
        comp = loc.get("competition", {})
        sal = loc.get("salary", {})
        rec = loc.get("recommendation_score", 0)

        _add_text_box(slide2, x + 0.15, y + 0.1, 3, 0.4,
                      name, font_size=11, bold=True, color=DOWNY_TEAL)
        _add_text_box(slide2, x + 3.5, y + 0.1, 2, 0.4,
                      f"Density: {td.get('score', 0)}", font_size=10, color=WHITE)
        _add_text_box(slide2, x + 5.5, y + 0.1, 2, 0.4,
                      f"Competition: {comp.get('score', 0)}", font_size=10, color=WHITE)
        _add_text_box(slide2, x + 7.5, y + 0.1, 2, 0.4,
                      f"Salary: ${sal.get('estimated', 0):,}", font_size=10, color=WHITE)

        rec_color = GREEN if rec >= 60 else (ORANGE if rec >= 40 else RED)
        _add_text_box(slide2, x + 9.8, y + 0.1, 1.5, 0.4,
                      f"Score: {rec}", font_size=10, bold=True, color=rec_color)
        y += 0.72

    # ── Slide 3: Salary Comparison ──
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide3, RGBColor(0x0F, 0x0F, 0x1E))
    _add_text_box(slide3, 0.8, 0.4, 8, 0.6,
                  "Salary Comparison", font_size=24, bold=True, color=WHITE)

    salary_data = analysis.get("salary_map", [])
    max_salary = max((s.get("estimated_salary", 1) for s in salary_data), default=1)

    y = 1.3
    for s in salary_data[:10]:
        name = s.get("display_name", s.get("location", ""))
        est = s.get("estimated_salary", 0)
        bar_w = max(0.3, (est / max_salary) * 8)

        _add_text_box(slide3, 0.8, y, 3, 0.35, name, font_size=10, color=WHITE)
        _add_shape(slide3, 4.0, y + 0.05, bar_w, 0.25, BLUE_VIOLET)
        _add_text_box(slide3, 4.0 + bar_w + 0.2, y, 2, 0.35,
                      f"${est:,}", font_size=10, bold=True, color=DOWNY_TEAL)
        y += 0.48

    # ── Slide 4: Optimal Locations ──
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide4, RGBColor(0x0F, 0x0F, 0x1E))
    _add_text_box(slide4, 0.8, 0.4, 8, 0.6,
                  "Optimal Hiring Locations", font_size=24, bold=True, color=WHITE)

    optimal = analysis.get("optimal_locations", [])
    y = 1.3
    card_w = 5.5
    for i, loc in enumerate(optimal[:6]):
        col = 0 if i % 2 == 0 else 1
        if i > 0 and col == 0:
            y += 2.0
        x = 0.8 if col == 0 else 7.0

        _add_shape(slide4, x, y, card_w, 1.8, CARD_BG)

        rank = loc.get("rank", i + 1)
        _add_text_box(slide4, x + 0.15, y + 0.1, 0.5, 0.35,
                      f"#{rank}", font_size=14, bold=True, color=DOWNY_TEAL)
        _add_text_box(slide4, x + 0.7, y + 0.1, 4.5, 0.35,
                      loc.get("location", ""), font_size=12, bold=True, color=WHITE)
        _add_text_box(slide4, x + 0.15, y + 0.55, 5, 0.3,
                      f"Score: {loc.get('composite_score', 0)} | Salary: ${loc.get('estimated_salary', 0):,} | CoL: {loc.get('coli', 100)}",
                      font_size=9, color=MUTED_TEXT)
        _add_text_box(slide4, x + 0.15, y + 0.95, 5, 0.7,
                      loc.get("reasoning", ""), font_size=9, color=WHITE)

    # Footer
    _add_text_box(slide4, 0.8, 6.8, 11, 0.4,
                  "Powered by Nova AI Suite | https://media-plan-generator.onrender.com",
                  font_size=8, color=MUTED_TEXT, alignment=PP_ALIGN.CENTER)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ═══════════════════════════════════════════════════════════════════════════════
# 11. ORCHESTRATOR -- run_heatmap_analysis
# ═══════════════════════════════════════════════════════════════════════════════

def run_heatmap_analysis(role: str, industry: str = "general_entry_level",
                         locations: Optional[List[str]] = None,
                         budget: float = 0, num_hires: int = 5) -> Dict[str, Any]:
    """Single orchestrator function for the Talent Supply Heat Map.

    Thread-safe, never crashes (wraps all stages in try/except).
    Runs analysis stages concurrently where possible.

    Args:
        role: Target job role/title
        industry: Canonical industry key
        locations: List of locations to analyze (auto-suggest if empty)
        budget: Total hiring budget in USD
        num_hires: Number of positions to fill

    Returns:
        Full analysis dict containing all sections.
    """
    start_time = time.time()

    result: Dict[str, Any] = {
        "status": "success",
        "role": role,
        "industry": industry,
        "industry_label": INDUSTRY_LABEL_MAP.get(industry, industry),
        "budget": budget,
        "num_hires": num_hires,
        "analysis_time_ms": 0,
        "errors": [],
    }

    # Validate inputs
    role = (role or "").strip()
    if not role:
        result["status"] = "error"
        result["errors"].append("Role is required")
        return result

    # Normalize industry -- map common names to canonical keys
    if industry not in INDUSTRY_LABEL_MAP:
        _industry_alias_map = {
            "technology": "tech_engineering", "tech": "tech_engineering",
            "engineering": "tech_engineering", "software": "tech_engineering",
            "it": "tech_engineering", "information technology": "tech_engineering",
            "healthcare": "healthcare_medical", "medical": "healthcare_medical",
            "health": "healthcare_medical", "nursing": "healthcare_medical",
            "finance": "finance_banking", "banking": "finance_banking",
            "financial": "finance_banking", "fintech": "finance_banking",
            "retail": "retail_consumer", "consumer": "retail_consumer",
            "ecommerce": "retail_consumer", "e-commerce": "retail_consumer",
            "logistics": "logistics_supply_chain", "supply chain": "logistics_supply_chain",
            "transportation": "logistics_supply_chain", "warehousing": "logistics_supply_chain",
            "hospitality": "hospitality_travel", "travel": "hospitality_travel",
            "hotel": "hospitality_travel", "restaurant": "hospitality_travel",
            "construction": "construction_real_estate", "real estate": "construction_real_estate",
            "education": "education", "teaching": "education",
            "aerospace": "aerospace_defense", "defense": "aerospace_defense",
            "pharma": "pharma_biotech", "biotech": "pharma_biotech",
            "pharmaceutical": "pharma_biotech",
            "energy": "energy_utilities", "utilities": "energy_utilities",
            "oil": "energy_utilities", "gas": "energy_utilities",
            "insurance": "insurance",
            "telecom": "telecommunications", "telecommunications": "telecommunications",
            "automotive": "automotive", "manufacturing": "automotive",
            "food": "food_beverage", "beverage": "food_beverage",
            "media": "media_entertainment", "entertainment": "media_entertainment",
            "legal": "legal_services", "law": "legal_services",
            "mental health": "mental_health", "behavioral": "mental_health",
            "blue collar": "blue_collar_trades", "trades": "blue_collar_trades",
            "skilled trades": "blue_collar_trades",
        }
        normalized = _industry_alias_map.get(industry.lower().strip(), None)
        if normalized:
            industry = normalized
        else:
            # Try partial match on label values
            _industry_lower = industry.lower().strip()
            matched = False
            for key, label in INDUSTRY_LABEL_MAP.items():
                if _industry_lower in label.lower() or label.lower() in _industry_lower:
                    industry = key
                    matched = True
                    break
            if not matched:
                industry = "general_entry_level"
        result["industry"] = industry
        result["industry_label"] = INDUSTRY_LABEL_MAP.get(industry, industry)

    # If no locations provided, auto-suggest top locations
    if not locations or (len(locations) == 1 and not locations[0].strip()):
        locations = _auto_suggest_locations(role, industry)
        result["auto_suggested"] = True
    else:
        locations = [l.strip() for l in locations if l and l.strip()]

    result["locations"] = locations

    if not locations:
        result["status"] = "error"
        result["errors"].append("No valid locations to analyze")
        return result

    # Run analysis stages concurrently
    with ThreadPoolExecutor(max_workers=5) as executor:
        f_comparison = executor.submit(
            _safe_call, compare_locations, role, locations, industry
        )
        f_salary = executor.submit(
            _safe_call, get_salary_map, role, locations, industry
        )
        f_optimal = executor.submit(
            _safe_call, find_optimal_locations, role, budget, num_hires, industry, 10
        )

        # Collect results
        try:
            result["comparison"] = f_comparison.result(timeout=30) or {}
        except Exception as exc:
            result["comparison"] = {}
            result["errors"].append(f"Comparison failed: {exc}")

        try:
            result["salary_map"] = f_salary.result(timeout=15) or []
        except Exception as exc:
            result["salary_map"] = []
            result["errors"].append(f"Salary map failed: {exc}")

        try:
            result["optimal_locations"] = f_optimal.result(timeout=30) or []
        except Exception as exc:
            result["optimal_locations"] = []
            result["errors"].append(f"Optimal locations failed: {exc}")

    # Calculate hiring difficulty per location
    try:
        difficulties = []
        for loc in locations:
            diff = calculate_hiring_difficulty(role, loc, industry)
            difficulties.append(diff)
        difficulties.sort(key=lambda x: x.get("difficulty_score", 0), reverse=True)
        result["difficulties"] = difficulties
    except Exception as exc:
        result["difficulties"] = []
        result["errors"].append(f"Hiring difficulty failed: {exc}")

    result["analysis_time_ms"] = int((time.time() - start_time) * 1000)

    if result["errors"]:
        result["status"] = "partial" if result.get("comparison") else "error"

    return result


def _auto_suggest_locations(role: str, industry: str) -> List[str]:
    """Auto-suggest top locations based on role and industry."""
    # Use industry-specific talent density to pick top metros
    density_map = INDUSTRY_TALENT_DENSITY.get(industry, {})

    if density_map:
        # Sort by density and return top 8
        sorted_locs = sorted(density_map.items(), key=lambda x: x[1], reverse=True)
        return [loc.replace("_", " ").title() for loc, _ in sorted_locs[:8]]

    # Generic top metros if no industry-specific data
    return [
        "New York", "San Francisco", "Chicago", "Dallas", "Austin",
        "Boston", "Seattle", "Denver",
    ]


# ═══════════════════════════════════════════════════════════════════════════════
# UTILITY HELPERS
# ═══════════════════════════════════════════════════════════════════════════════

def _safe_call(fn, *args, **kwargs):
    """Call a function, returning None on any exception."""
    try:
        return fn(*args, **kwargs)
    except Exception as exc:
        logger.warning("_safe_call(%s) failed: %s", fn.__name__, exc)
        return None


def get_industry_options() -> List[Dict[str, str]]:
    """Return industry options for the frontend dropdown."""
    return [
        {"value": key, "label": label}
        for key, label in INDUSTRY_LABEL_MAP.items()
    ]
