"""
skill_target.py -- Skill-Based Targeting Engine

Uses O*NET occupation data + collar_intelligence to help recruiters target
candidates by SKILLS rather than just job titles. Provides skill-to-occupation
matching, demand trends, channel recommendations, salary benchmarks,
geographic hotspots, adjacent skills, and exportable Excel/PPT reports.

Thread-safe. All public functions catch exceptions and return safe fallbacks.
"""

from __future__ import annotations

import io
import logging
import re
from concurrent.futures import ThreadPoolExecutor, as_completed
from typing import Any, Dict, List, Optional, Tuple

logger = logging.getLogger(__name__)

# -- Lazy / optional imports --

try:
    from collar_intelligence import (
        classify_collar,
        get_blended_allocation,
        get_collar_comparison,
        analyze_skills_gap,
        COLLAR_STRATEGY,
        ROLE_SKILLS_MAP,
        SKILL_SCARCITY,
    )

    _HAS_COLLAR = True
except ImportError:
    _HAS_COLLAR = False
    COLLAR_STRATEGY = {}
    ROLE_SKILLS_MAP = {}
    SKILL_SCARCITY = {}

try:
    from research import (
        get_market_trends,
        get_location_info,
        get_location_boards,
        get_media_platform_audiences,
    )

    _HAS_RESEARCH = True
except ImportError:
    _HAS_RESEARCH = False

try:
    from api_enrichment import fetch_salary_data, fetch_location_demographics

    _HAS_API = True
except ImportError:
    _HAS_API = False

try:
    from shared_utils import INDUSTRY_LABEL_MAP, parse_budget

    _HAS_UTILS = True
except ImportError:
    _HAS_UTILS = False
    INDUSTRY_LABEL_MAP = {}


# ===========================================================================
# 1. SKILLS-TO-OCCUPATIONS MAPPING  (O*NET-aligned, 60+ skills)
# ===========================================================================

SKILLS_TO_OCCUPATIONS: Dict[str, List[Dict[str, str]]] = {
    "Python": [
        {"soc": "15-1252", "title": "Software Developer", "zone": "4"},
        {"soc": "15-2051", "title": "Data Scientist", "zone": "5"},
        {"soc": "15-1211", "title": "Computer Systems Analyst", "zone": "4"},
        {"soc": "15-1299", "title": "Computer Occupations, All Other", "zone": "4"},
    ],
    "JavaScript": [
        {"soc": "15-1254", "title": "Web Developer", "zone": "3"},
        {"soc": "15-1252", "title": "Software Developer", "zone": "4"},
        {"soc": "15-1255", "title": "Web and Digital Interface Designer", "zone": "4"},
    ],
    "SQL": [
        {"soc": "15-1243", "title": "Database Administrator", "zone": "4"},
        {"soc": "15-2051", "title": "Data Scientist", "zone": "5"},
        {"soc": "15-1252", "title": "Software Developer", "zone": "4"},
        {"soc": "13-1111", "title": "Management Analyst", "zone": "4"},
    ],
    "Cloud (AWS/GCP/Azure)": [
        {
            "soc": "15-1244",
            "title": "Network and Computer Systems Administrator",
            "zone": "4",
        },
        {"soc": "15-1252", "title": "Software Developer", "zone": "4"},
        {"soc": "15-1212", "title": "Information Security Analyst", "zone": "4"},
    ],
    "Machine Learning": [
        {"soc": "15-2051", "title": "Data Scientist", "zone": "5"},
        {"soc": "15-1252", "title": "Software Developer", "zone": "4"},
        {"soc": "15-2041", "title": "Statistician", "zone": "5"},
    ],
    "Cybersecurity": [
        {"soc": "15-1212", "title": "Information Security Analyst", "zone": "4"},
        {
            "soc": "15-1244",
            "title": "Network and Computer Systems Administrator",
            "zone": "4",
        },
    ],
    "DevOps": [
        {
            "soc": "15-1244",
            "title": "Network and Computer Systems Administrator",
            "zone": "4",
        },
        {"soc": "15-1252", "title": "Software Developer", "zone": "4"},
    ],
    "Kubernetes": [
        {
            "soc": "15-1244",
            "title": "Network and Computer Systems Administrator",
            "zone": "4",
        },
        {"soc": "15-1252", "title": "Software Developer", "zone": "4"},
    ],
    "React": [
        {"soc": "15-1254", "title": "Web Developer", "zone": "3"},
        {"soc": "15-1252", "title": "Software Developer", "zone": "4"},
    ],
    "TypeScript": [
        {"soc": "15-1254", "title": "Web Developer", "zone": "3"},
        {"soc": "15-1252", "title": "Software Developer", "zone": "4"},
    ],
    "Java": [
        {"soc": "15-1252", "title": "Software Developer", "zone": "4"},
        {"soc": "15-1299", "title": "Computer Occupations, All Other", "zone": "4"},
    ],
    "C++": [
        {"soc": "15-1252", "title": "Software Developer", "zone": "4"},
        {"soc": "17-2061", "title": "Computer Hardware Engineer", "zone": "4"},
    ],
    "Rust": [{"soc": "15-1252", "title": "Software Developer", "zone": "4"}],
    "Go": [
        {"soc": "15-1252", "title": "Software Developer", "zone": "4"},
        {
            "soc": "15-1244",
            "title": "Network and Computer Systems Administrator",
            "zone": "4",
        },
    ],
    "Data Visualization": [
        {"soc": "15-2051", "title": "Data Scientist", "zone": "5"},
        {"soc": "15-1257", "title": "Web and Digital Interface Designer", "zone": "4"},
    ],
    "TensorFlow/PyTorch": [
        {"soc": "15-2051", "title": "Data Scientist", "zone": "5"},
        {"soc": "15-1252", "title": "Software Developer", "zone": "4"},
    ],
    "System Design": [
        {"soc": "15-1252", "title": "Software Developer", "zone": "4"},
        {"soc": "15-1241", "title": "Computer Network Architect", "zone": "4"},
    ],
    "Git": [
        {"soc": "15-1252", "title": "Software Developer", "zone": "4"},
        {"soc": "15-1254", "title": "Web Developer", "zone": "3"},
    ],
    "Patient Assessment": [
        {"soc": "29-1141", "title": "Registered Nurse", "zone": "3"},
        {"soc": "29-1071", "title": "Physician Assistant", "zone": "5"},
        {"soc": "29-1216", "title": "General Internal Medicine Physician", "zone": "5"},
    ],
    "BLS/ACLS": [
        {"soc": "29-1141", "title": "Registered Nurse", "zone": "3"},
        {"soc": "29-2041", "title": "Emergency Medical Technician", "zone": "2"},
    ],
    "Electronic Health Records": [
        {"soc": "29-1141", "title": "Registered Nurse", "zone": "3"},
        {"soc": "29-2072", "title": "Medical Records Specialist", "zone": "3"},
    ],
    "Medication Administration": [
        {"soc": "29-1141", "title": "Registered Nurse", "zone": "3"},
        {"soc": "29-2052", "title": "Pharmacy Technician", "zone": "2"},
    ],
    "Clinical Diagnosis": [
        {"soc": "29-1216", "title": "General Internal Medicine Physician", "zone": "5"},
        {"soc": "29-1071", "title": "Physician Assistant", "zone": "5"},
    ],
    "Therapeutic Techniques": [
        {"soc": "21-1014", "title": "Mental Health Counselor", "zone": "5"},
        {"soc": "29-1223", "title": "Psychiatrist", "zone": "5"},
        {"soc": "29-1122", "title": "Occupational Therapist", "zone": "5"},
    ],
    "Phlebotomy": [
        {"soc": "31-9097", "title": "Phlebotomist", "zone": "2"},
        {
            "soc": "29-2012",
            "title": "Medical and Clinical Laboratory Technician",
            "zone": "3",
        },
    ],
    "Medical Coding": [
        {"soc": "29-2072", "title": "Medical Records Specialist", "zone": "3"},
        {"soc": "43-6013", "title": "Medical Secretary", "zone": "3"},
    ],
    "CDL License": [
        {
            "soc": "53-3032",
            "title": "Heavy and Tractor-Trailer Truck Driver",
            "zone": "2",
        },
        {"soc": "53-3033", "title": "Light Truck Driver", "zone": "2"},
    ],
    "Welding": [
        {"soc": "51-4121", "title": "Welder, Cutter, Solderer", "zone": "3"},
        {"soc": "47-2152", "title": "Plumber, Pipefitter", "zone": "3"},
    ],
    "Forklift Operation": [
        {
            "soc": "53-7051",
            "title": "Industrial Truck and Tractor Operator",
            "zone": "1",
        },
        {"soc": "53-7062", "title": "Laborer and Material Mover", "zone": "1"},
    ],
    "Blueprint Reading": [
        {"soc": "47-2111", "title": "Electrician", "zone": "3"},
        {"soc": "47-2152", "title": "Plumber, Pipefitter", "zone": "3"},
        {"soc": "47-2061", "title": "Construction Laborer", "zone": "1"},
    ],
    "HVAC": [{"soc": "49-9021", "title": "HVAC Mechanic and Installer", "zone": "3"}],
    "Electrical Troubleshooting": [
        {"soc": "47-2111", "title": "Electrician", "zone": "3"},
        {
            "soc": "49-2098",
            "title": "Security and Fire Alarm Systems Installer",
            "zone": "3",
        },
    ],
    "CNC Machining": [
        {"soc": "51-4011", "title": "CNC Tool Operator", "zone": "3"},
        {"soc": "51-4041", "title": "Machinist", "zone": "3"},
    ],
    "Plumbing": [{"soc": "47-2152", "title": "Plumber, Pipefitter", "zone": "3"}],
    "Carpentry": [{"soc": "47-2031", "title": "Carpenter", "zone": "2"}],
    "Heavy Equipment Operation": [
        {"soc": "47-2073", "title": "Operating Engineer", "zone": "2"},
        {"soc": "53-7032", "title": "Excavating Operator", "zone": "2"},
    ],
    "Project Management": [
        {"soc": "11-9199", "title": "Manager, All Other", "zone": "4"},
        {"soc": "13-1082", "title": "Project Management Specialist", "zone": "4"},
    ],
    "Data Analysis": [
        {"soc": "15-2051", "title": "Data Scientist", "zone": "5"},
        {"soc": "13-1111", "title": "Management Analyst", "zone": "4"},
    ],
    "Financial Modeling": [
        {"soc": "13-2051", "title": "Financial Analyst", "zone": "4"},
        {"soc": "13-2054", "title": "Financial Risk Specialist", "zone": "4"},
    ],
    "SEO/SEM": [
        {"soc": "13-1161", "title": "Market Research Analyst", "zone": "4"},
        {"soc": "27-3031", "title": "Public Relations Specialist", "zone": "4"},
    ],
    "Digital Marketing": [
        {"soc": "13-1161", "title": "Market Research Analyst", "zone": "4"},
        {"soc": "11-2021", "title": "Marketing Manager", "zone": "4"},
    ],
    "Salesforce": [
        {"soc": "15-1252", "title": "Software Developer", "zone": "4"},
        {"soc": "15-1211", "title": "Computer Systems Analyst", "zone": "4"},
    ],
    "SAP": [
        {"soc": "15-1211", "title": "Computer Systems Analyst", "zone": "4"},
        {"soc": "15-1252", "title": "Software Developer", "zone": "4"},
    ],
    "Accounting (GAAP)": [
        {"soc": "13-2011", "title": "Accountant and Auditor", "zone": "4"},
        {"soc": "13-2082", "title": "Tax Preparer", "zone": "3"},
    ],
    "CRM Software": [
        {"soc": "41-3091", "title": "Sales Representative", "zone": "3"},
        {"soc": "11-2022", "title": "Sales Manager", "zone": "4"},
    ],
    "Negotiation": [
        {"soc": "41-3091", "title": "Sales Representative", "zone": "3"},
        {"soc": "13-1075", "title": "Labor Relations Specialist", "zone": "4"},
    ],
    "UX Design": [
        {"soc": "15-1255", "title": "Web and Digital Interface Designer", "zone": "4"},
        {"soc": "27-1024", "title": "Graphic Designer", "zone": "3"},
    ],
    "Figma/Sketch": [
        {"soc": "15-1255", "title": "Web and Digital Interface Designer", "zone": "4"},
        {"soc": "27-1024", "title": "Graphic Designer", "zone": "3"},
    ],
    "AutoCAD": [
        {"soc": "17-3011", "title": "Architectural and Civil Drafter", "zone": "3"},
        {"soc": "17-2051", "title": "Civil Engineer", "zone": "4"},
    ],
    "Six Sigma": [
        {"soc": "17-2112", "title": "Industrial Engineer", "zone": "4"},
        {"soc": "13-1082", "title": "Project Management Specialist", "zone": "4"},
    ],
    "Supply Chain Management": [
        {"soc": "13-1081", "title": "Logistician", "zone": "4"},
        {
            "soc": "11-3071",
            "title": "Transportation, Storage, Distribution Manager",
            "zone": "4",
        },
    ],
    "Customer Service": [
        {"soc": "43-4051", "title": "Customer Service Representative", "zone": "2"},
        {"soc": "41-2031", "title": "Retail Salesperson", "zone": "2"},
    ],
    "Food Safety (ServSafe)": [
        {
            "soc": "35-1012",
            "title": "First-Line Supervisor of Food Preparation",
            "zone": "2",
        },
        {"soc": "35-2014", "title": "Cook, Restaurant", "zone": "2"},
    ],
    "POS Systems": [
        {"soc": "41-2011", "title": "Cashier", "zone": "1"},
        {"soc": "41-2031", "title": "Retail Salesperson", "zone": "2"},
    ],
    "Curriculum Development": [
        {"soc": "25-1000", "title": "Postsecondary Teacher", "zone": "5"},
        {"soc": "25-2021", "title": "Elementary School Teacher", "zone": "4"},
        {"soc": "25-9031", "title": "Instructional Coordinator", "zone": "5"},
    ],
    "Legal Research": [
        {"soc": "23-1011", "title": "Lawyer", "zone": "5"},
        {"soc": "23-2011", "title": "Paralegal and Legal Assistant", "zone": "3"},
    ],
    "Contract Negotiation": [
        {"soc": "23-1011", "title": "Lawyer", "zone": "5"},
        {"soc": "13-1075", "title": "Labor Relations Specialist", "zone": "4"},
    ],
}


# ===========================================================================
# 2. ADJACENT SKILLS GRAPH
# ===========================================================================

ADJACENT_SKILLS: Dict[str, List[Dict[str, Any]]] = {
    "Python": [
        {"skill": "SQL", "relevance": 0.85},
        {"skill": "Machine Learning", "relevance": 0.75},
        {"skill": "Data Analysis", "relevance": 0.80},
        {"skill": "Git", "relevance": 0.70},
        {"skill": "Cloud (AWS/GCP/Azure)", "relevance": 0.65},
        {"skill": "TensorFlow/PyTorch", "relevance": 0.60},
    ],
    "JavaScript": [
        {"skill": "React", "relevance": 0.90},
        {"skill": "TypeScript", "relevance": 0.88},
        {"skill": "Git", "relevance": 0.75},
        {"skill": "SQL", "relevance": 0.55},
        {"skill": "UX Design", "relevance": 0.50},
    ],
    "SQL": [
        {"skill": "Python", "relevance": 0.80},
        {"skill": "Data Analysis", "relevance": 0.85},
        {"skill": "Data Visualization", "relevance": 0.70},
        {"skill": "Cloud (AWS/GCP/Azure)", "relevance": 0.55},
    ],
    "Cloud (AWS/GCP/Azure)": [
        {"skill": "Kubernetes", "relevance": 0.85},
        {"skill": "DevOps", "relevance": 0.88},
        {"skill": "Python", "relevance": 0.65},
        {"skill": "Cybersecurity", "relevance": 0.60},
        {"skill": "System Design", "relevance": 0.75},
    ],
    "Machine Learning": [
        {"skill": "Python", "relevance": 0.90},
        {"skill": "TensorFlow/PyTorch", "relevance": 0.92},
        {"skill": "Data Visualization", "relevance": 0.65},
        {"skill": "SQL", "relevance": 0.60},
        {"skill": "Cloud (AWS/GCP/Azure)", "relevance": 0.55},
    ],
    "Cybersecurity": [
        {"skill": "Cloud (AWS/GCP/Azure)", "relevance": 0.75},
        {"skill": "Python", "relevance": 0.60},
        {"skill": "DevOps", "relevance": 0.55},
        {"skill": "System Design", "relevance": 0.50},
    ],
    "DevOps": [
        {"skill": "Kubernetes", "relevance": 0.90},
        {"skill": "Cloud (AWS/GCP/Azure)", "relevance": 0.88},
        {"skill": "Git", "relevance": 0.80},
        {"skill": "Python", "relevance": 0.60},
    ],
    "Kubernetes": [
        {"skill": "DevOps", "relevance": 0.90},
        {"skill": "Cloud (AWS/GCP/Azure)", "relevance": 0.88},
        {"skill": "Go", "relevance": 0.50},
    ],
    "React": [
        {"skill": "JavaScript", "relevance": 0.95},
        {"skill": "TypeScript", "relevance": 0.85},
        {"skill": "UX Design", "relevance": 0.55},
        {"skill": "Git", "relevance": 0.70},
    ],
    "TypeScript": [
        {"skill": "JavaScript", "relevance": 0.95},
        {"skill": "React", "relevance": 0.80},
        {"skill": "Git", "relevance": 0.65},
    ],
    "Java": [
        {"skill": "SQL", "relevance": 0.75},
        {"skill": "System Design", "relevance": 0.70},
        {"skill": "Cloud (AWS/GCP/Azure)", "relevance": 0.60},
        {"skill": "Git", "relevance": 0.70},
    ],
    "C++": [
        {"skill": "System Design", "relevance": 0.80},
        {"skill": "Rust", "relevance": 0.55},
        {"skill": "Python", "relevance": 0.50},
    ],
    "Rust": [
        {"skill": "C++", "relevance": 0.70},
        {"skill": "System Design", "relevance": 0.75},
        {"skill": "Go", "relevance": 0.50},
    ],
    "Go": [
        {"skill": "Kubernetes", "relevance": 0.65},
        {"skill": "DevOps", "relevance": 0.55},
        {"skill": "Cloud (AWS/GCP/Azure)", "relevance": 0.60},
    ],
    "Data Visualization": [
        {"skill": "Data Analysis", "relevance": 0.90},
        {"skill": "SQL", "relevance": 0.70},
        {"skill": "Python", "relevance": 0.65},
    ],
    "TensorFlow/PyTorch": [
        {"skill": "Machine Learning", "relevance": 0.95},
        {"skill": "Python", "relevance": 0.90},
        {"skill": "Data Visualization", "relevance": 0.50},
    ],
    "System Design": [
        {"skill": "Cloud (AWS/GCP/Azure)", "relevance": 0.80},
        {"skill": "DevOps", "relevance": 0.60},
        {"skill": "Kubernetes", "relevance": 0.55},
    ],
    "Salesforce": [
        {"skill": "CRM Software", "relevance": 0.90},
        {"skill": "SQL", "relevance": 0.55},
        {"skill": "Project Management", "relevance": 0.45},
    ],
    "SAP": [
        {"skill": "Supply Chain Management", "relevance": 0.70},
        {"skill": "SQL", "relevance": 0.55},
        {"skill": "Financial Modeling", "relevance": 0.50},
    ],
    "Patient Assessment": [
        {"skill": "BLS/ACLS", "relevance": 0.85},
        {"skill": "Electronic Health Records", "relevance": 0.80},
        {"skill": "Medication Administration", "relevance": 0.75},
        {"skill": "Clinical Diagnosis", "relevance": 0.65},
    ],
    "BLS/ACLS": [
        {"skill": "Patient Assessment", "relevance": 0.85},
        {"skill": "Medication Administration", "relevance": 0.70},
    ],
    "Electronic Health Records": [
        {"skill": "Patient Assessment", "relevance": 0.75},
        {"skill": "Medical Coding", "relevance": 0.70},
    ],
    "Therapeutic Techniques": [
        {"skill": "Patient Assessment", "relevance": 0.70},
        {"skill": "Clinical Diagnosis", "relevance": 0.65},
    ],
    "Phlebotomy": [
        {"skill": "Patient Assessment", "relevance": 0.55},
        {"skill": "BLS/ACLS", "relevance": 0.45},
    ],
    "Medical Coding": [
        {"skill": "Electronic Health Records", "relevance": 0.80},
        {"skill": "Accounting (GAAP)", "relevance": 0.30},
    ],
    "CDL License": [
        {"skill": "Forklift Operation", "relevance": 0.50},
        {"skill": "Heavy Equipment Operation", "relevance": 0.45},
    ],
    "Welding": [
        {"skill": "Blueprint Reading", "relevance": 0.75},
        {"skill": "CNC Machining", "relevance": 0.55},
        {"skill": "Plumbing", "relevance": 0.45},
    ],
    "Blueprint Reading": [
        {"skill": "Welding", "relevance": 0.65},
        {"skill": "Electrical Troubleshooting", "relevance": 0.70},
        {"skill": "Plumbing", "relevance": 0.65},
        {"skill": "Carpentry", "relevance": 0.60},
        {"skill": "AutoCAD", "relevance": 0.55},
    ],
    "HVAC": [
        {"skill": "Electrical Troubleshooting", "relevance": 0.75},
        {"skill": "Plumbing", "relevance": 0.55},
        {"skill": "Blueprint Reading", "relevance": 0.60},
    ],
    "CNC Machining": [
        {"skill": "Blueprint Reading", "relevance": 0.80},
        {"skill": "Welding", "relevance": 0.55},
        {"skill": "Six Sigma", "relevance": 0.40},
    ],
    "Electrical Troubleshooting": [
        {"skill": "Blueprint Reading", "relevance": 0.75},
        {"skill": "HVAC", "relevance": 0.60},
    ],
    "Heavy Equipment Operation": [
        {"skill": "CDL License", "relevance": 0.55},
        {"skill": "Carpentry", "relevance": 0.35},
    ],
    "Forklift Operation": [
        {"skill": "CDL License", "relevance": 0.40},
        {"skill": "Supply Chain Management", "relevance": 0.35},
    ],
    "Carpentry": [
        {"skill": "Blueprint Reading", "relevance": 0.75},
        {"skill": "Heavy Equipment Operation", "relevance": 0.35},
    ],
    "Plumbing": [
        {"skill": "Blueprint Reading", "relevance": 0.70},
        {"skill": "Welding", "relevance": 0.55},
        {"skill": "HVAC", "relevance": 0.50},
    ],
    "Project Management": [
        {"skill": "Six Sigma", "relevance": 0.55},
        {"skill": "Data Analysis", "relevance": 0.50},
        {"skill": "Negotiation", "relevance": 0.45},
        {"skill": "Supply Chain Management", "relevance": 0.40},
    ],
    "Data Analysis": [
        {"skill": "SQL", "relevance": 0.85},
        {"skill": "Python", "relevance": 0.75},
        {"skill": "Data Visualization", "relevance": 0.88},
    ],
    "Financial Modeling": [
        {"skill": "Accounting (GAAP)", "relevance": 0.80},
        {"skill": "Data Analysis", "relevance": 0.65},
        {"skill": "SQL", "relevance": 0.50},
    ],
    "SEO/SEM": [
        {"skill": "Digital Marketing", "relevance": 0.90},
        {"skill": "Data Analysis", "relevance": 0.55},
    ],
    "Digital Marketing": [
        {"skill": "SEO/SEM", "relevance": 0.90},
        {"skill": "Data Analysis", "relevance": 0.55},
        {"skill": "CRM Software", "relevance": 0.50},
    ],
    "CRM Software": [
        {"skill": "Salesforce", "relevance": 0.85},
        {"skill": "Digital Marketing", "relevance": 0.50},
        {"skill": "Negotiation", "relevance": 0.45},
    ],
    "Negotiation": [
        {"skill": "CRM Software", "relevance": 0.50},
        {"skill": "Project Management", "relevance": 0.45},
        {"skill": "Contract Negotiation", "relevance": 0.85},
    ],
    "UX Design": [
        {"skill": "Figma/Sketch", "relevance": 0.90},
        {"skill": "JavaScript", "relevance": 0.50},
        {"skill": "React", "relevance": 0.40},
    ],
    "Figma/Sketch": [
        {"skill": "UX Design", "relevance": 0.90},
        {"skill": "React", "relevance": 0.40},
    ],
    "AutoCAD": [
        {"skill": "Blueprint Reading", "relevance": 0.80},
        {"skill": "Six Sigma", "relevance": 0.30},
    ],
    "Six Sigma": [
        {"skill": "Project Management", "relevance": 0.65},
        {"skill": "Data Analysis", "relevance": 0.55},
        {"skill": "Supply Chain Management", "relevance": 0.60},
    ],
    "Supply Chain Management": [
        {"skill": "Project Management", "relevance": 0.60},
        {"skill": "SAP", "relevance": 0.65},
        {"skill": "Six Sigma", "relevance": 0.55},
    ],
    "Accounting (GAAP)": [
        {"skill": "Financial Modeling", "relevance": 0.80},
        {"skill": "SAP", "relevance": 0.45},
    ],
    "Customer Service": [
        {"skill": "POS Systems", "relevance": 0.60},
        {"skill": "CRM Software", "relevance": 0.45},
        {"skill": "Negotiation", "relevance": 0.35},
    ],
    "Food Safety (ServSafe)": [{"skill": "Customer Service", "relevance": 0.50}],
    "Curriculum Development": [
        {"skill": "Data Analysis", "relevance": 0.35},
        {"skill": "Project Management", "relevance": 0.40},
    ],
    "Legal Research": [
        {"skill": "Contract Negotiation", "relevance": 0.80},
        {"skill": "Data Analysis", "relevance": 0.40},
    ],
    "Contract Negotiation": [
        {"skill": "Legal Research", "relevance": 0.80},
        {"skill": "Negotiation", "relevance": 0.85},
    ],
}


# ===========================================================================
# 3. SKILL DEMAND INDICATORS
# ===========================================================================

_SKILL_DEMAND: Dict[str, Dict[str, Any]] = {
    "Python": {"growth": "high", "yoy_pct": 12, "pool": "large", "shortage": False},
    "JavaScript": {"growth": "high", "yoy_pct": 8, "pool": "large", "shortage": False},
    "SQL": {"growth": "stable", "yoy_pct": 4, "pool": "large", "shortage": False},
    "Cloud (AWS/GCP/Azure)": {
        "growth": "very_high",
        "yoy_pct": 22,
        "pool": "medium",
        "shortage": True,
    },
    "Machine Learning": {
        "growth": "very_high",
        "yoy_pct": 28,
        "pool": "small",
        "shortage": True,
    },
    "Cybersecurity": {
        "growth": "very_high",
        "yoy_pct": 25,
        "pool": "small",
        "shortage": True,
    },
    "DevOps": {"growth": "high", "yoy_pct": 18, "pool": "medium", "shortage": True},
    "Kubernetes": {
        "growth": "very_high",
        "yoy_pct": 30,
        "pool": "small",
        "shortage": True,
    },
    "React": {"growth": "high", "yoy_pct": 10, "pool": "large", "shortage": False},
    "TypeScript": {
        "growth": "high",
        "yoy_pct": 20,
        "pool": "medium",
        "shortage": False,
    },
    "Java": {"growth": "stable", "yoy_pct": 3, "pool": "large", "shortage": False},
    "C++": {"growth": "stable", "yoy_pct": 2, "pool": "medium", "shortage": False},
    "Rust": {"growth": "high", "yoy_pct": 35, "pool": "small", "shortage": True},
    "Go": {"growth": "high", "yoy_pct": 18, "pool": "small", "shortage": True},
    "Data Visualization": {
        "growth": "high",
        "yoy_pct": 10,
        "pool": "medium",
        "shortage": False,
    },
    "TensorFlow/PyTorch": {
        "growth": "very_high",
        "yoy_pct": 25,
        "pool": "small",
        "shortage": True,
    },
    "System Design": {
        "growth": "high",
        "yoy_pct": 12,
        "pool": "medium",
        "shortage": True,
    },
    "Git": {"growth": "stable", "yoy_pct": 3, "pool": "large", "shortage": False},
    "Salesforce": {"growth": "high", "yoy_pct": 14, "pool": "medium", "shortage": True},
    "SAP": {"growth": "stable", "yoy_pct": 5, "pool": "medium", "shortage": True},
    "Patient Assessment": {
        "growth": "high",
        "yoy_pct": 8,
        "pool": "medium",
        "shortage": True,
    },
    "BLS/ACLS": {"growth": "stable", "yoy_pct": 5, "pool": "medium", "shortage": False},
    "Electronic Health Records": {
        "growth": "high",
        "yoy_pct": 10,
        "pool": "medium",
        "shortage": False,
    },
    "Medication Administration": {
        "growth": "stable",
        "yoy_pct": 5,
        "pool": "medium",
        "shortage": True,
    },
    "Clinical Diagnosis": {
        "growth": "stable",
        "yoy_pct": 4,
        "pool": "small",
        "shortage": True,
    },
    "Therapeutic Techniques": {
        "growth": "high",
        "yoy_pct": 12,
        "pool": "small",
        "shortage": True,
    },
    "Phlebotomy": {
        "growth": "stable",
        "yoy_pct": 6,
        "pool": "medium",
        "shortage": False,
    },
    "Medical Coding": {
        "growth": "high",
        "yoy_pct": 9,
        "pool": "medium",
        "shortage": False,
    },
    "CDL License": {"growth": "high", "yoy_pct": 8, "pool": "medium", "shortage": True},
    "Welding": {"growth": "high", "yoy_pct": 10, "pool": "medium", "shortage": True},
    "Forklift Operation": {
        "growth": "stable",
        "yoy_pct": 4,
        "pool": "large",
        "shortage": False,
    },
    "Blueprint Reading": {
        "growth": "stable",
        "yoy_pct": 3,
        "pool": "medium",
        "shortage": False,
    },
    "HVAC": {"growth": "high", "yoy_pct": 12, "pool": "small", "shortage": True},
    "CNC Machining": {
        "growth": "stable",
        "yoy_pct": 5,
        "pool": "small",
        "shortage": True,
    },
    "Electrical Troubleshooting": {
        "growth": "high",
        "yoy_pct": 8,
        "pool": "medium",
        "shortage": True,
    },
    "Heavy Equipment Operation": {
        "growth": "stable",
        "yoy_pct": 4,
        "pool": "medium",
        "shortage": False,
    },
    "Carpentry": {
        "growth": "stable",
        "yoy_pct": 3,
        "pool": "medium",
        "shortage": False,
    },
    "Plumbing": {"growth": "high", "yoy_pct": 8, "pool": "medium", "shortage": True},
    "Project Management": {
        "growth": "stable",
        "yoy_pct": 5,
        "pool": "large",
        "shortage": False,
    },
    "Data Analysis": {
        "growth": "high",
        "yoy_pct": 14,
        "pool": "large",
        "shortage": False,
    },
    "Financial Modeling": {
        "growth": "stable",
        "yoy_pct": 5,
        "pool": "medium",
        "shortage": False,
    },
    "SEO/SEM": {"growth": "high", "yoy_pct": 8, "pool": "medium", "shortage": False},
    "Digital Marketing": {
        "growth": "high",
        "yoy_pct": 10,
        "pool": "large",
        "shortage": False,
    },
    "CRM Software": {
        "growth": "stable",
        "yoy_pct": 4,
        "pool": "large",
        "shortage": False,
    },
    "Negotiation": {
        "growth": "stable",
        "yoy_pct": 2,
        "pool": "large",
        "shortage": False,
    },
    "UX Design": {"growth": "high", "yoy_pct": 12, "pool": "medium", "shortage": False},
    "Figma/Sketch": {
        "growth": "high",
        "yoy_pct": 15,
        "pool": "medium",
        "shortage": False,
    },
    "AutoCAD": {"growth": "stable", "yoy_pct": 3, "pool": "medium", "shortage": False},
    "Six Sigma": {
        "growth": "stable",
        "yoy_pct": 2,
        "pool": "medium",
        "shortage": False,
    },
    "Supply Chain Management": {
        "growth": "high",
        "yoy_pct": 10,
        "pool": "medium",
        "shortage": True,
    },
    "Accounting (GAAP)": {
        "growth": "stable",
        "yoy_pct": 3,
        "pool": "large",
        "shortage": False,
    },
    "Customer Service": {
        "growth": "stable",
        "yoy_pct": 1,
        "pool": "large",
        "shortage": False,
    },
    "Food Safety (ServSafe)": {
        "growth": "stable",
        "yoy_pct": 3,
        "pool": "large",
        "shortage": False,
    },
    "Curriculum Development": {
        "growth": "stable",
        "yoy_pct": 4,
        "pool": "medium",
        "shortage": False,
    },
    "Legal Research": {
        "growth": "stable",
        "yoy_pct": 2,
        "pool": "medium",
        "shortage": False,
    },
    "Contract Negotiation": {
        "growth": "stable",
        "yoy_pct": 3,
        "pool": "medium",
        "shortage": False,
    },
    "POS Systems": {
        "growth": "stable",
        "yoy_pct": 1,
        "pool": "large",
        "shortage": False,
    },
}


# ===========================================================================
# 4. GEOGRAPHIC HOTSPOT DATA
# ===========================================================================

_SKILL_HOTSPOTS: Dict[str, List[Dict[str, Any]]] = {
    "Python": [
        {
            "metro": "San Francisco-Oakland, CA",
            "concentration": 0.92,
            "employers": "Google, Meta, Stripe",
        },
        {
            "metro": "Seattle-Tacoma, WA",
            "concentration": 0.88,
            "employers": "Amazon, Microsoft",
        },
        {
            "metro": "New York-Newark, NY",
            "concentration": 0.82,
            "employers": "JP Morgan, Bloomberg",
        },
        {
            "metro": "Austin-Round Rock, TX",
            "concentration": 0.78,
            "employers": "Dell, Oracle, Tesla",
        },
        {
            "metro": "Boston-Cambridge, MA",
            "concentration": 0.76,
            "employers": "HubSpot, Wayfair",
        },
    ],
    "JavaScript": [
        {
            "metro": "San Francisco-Oakland, CA",
            "concentration": 0.90,
            "employers": "Meta, Vercel, Stripe",
        },
        {
            "metro": "New York-Newark, NY",
            "concentration": 0.85,
            "employers": "Squarespace, Datadog",
        },
        {
            "metro": "Seattle-Tacoma, WA",
            "concentration": 0.82,
            "employers": "Amazon, Microsoft",
        },
    ],
    "Cloud (AWS/GCP/Azure)": [
        {
            "metro": "Seattle-Tacoma, WA",
            "concentration": 0.95,
            "employers": "Amazon AWS, Microsoft Azure",
        },
        {
            "metro": "Northern Virginia (Ashburn)",
            "concentration": 0.90,
            "employers": "AWS, Rackspace",
        },
        {
            "metro": "San Francisco-Oakland, CA",
            "concentration": 0.88,
            "employers": "Google Cloud, Salesforce",
        },
        {
            "metro": "Dallas-Fort Worth, TX",
            "concentration": 0.72,
            "employers": "AT&T, data centers",
        },
    ],
    "Machine Learning": [
        {
            "metro": "San Francisco-Oakland, CA",
            "concentration": 0.95,
            "employers": "OpenAI, Google DeepMind, Anthropic",
        },
        {
            "metro": "Seattle-Tacoma, WA",
            "concentration": 0.85,
            "employers": "Amazon, Allen AI",
        },
        {
            "metro": "Boston-Cambridge, MA",
            "concentration": 0.80,
            "employers": "MIT, Moderna",
        },
        {
            "metro": "New York-Newark, NY",
            "concentration": 0.75,
            "employers": "Bloomberg, Two Sigma",
        },
    ],
    "Cybersecurity": [
        {
            "metro": "Washington-Arlington, DC",
            "concentration": 0.95,
            "employers": "NSA, CrowdStrike, Booz Allen",
        },
        {
            "metro": "Northern Virginia (Ashburn)",
            "concentration": 0.90,
            "employers": "Mandiant, Palo Alto Networks",
        },
        {
            "metro": "San Francisco-Oakland, CA",
            "concentration": 0.78,
            "employers": "Cloudflare, Zscaler",
        },
    ],
    "CDL License": [
        {
            "metro": "Dallas-Fort Worth, TX",
            "concentration": 0.88,
            "employers": "FedEx, UPS, Werner",
        },
        {
            "metro": "Atlanta-Sandy Springs, GA",
            "concentration": 0.85,
            "employers": "UPS, Ryder",
        },
        {
            "metro": "Chicago-Naperville, IL",
            "concentration": 0.82,
            "employers": "Schneider, XPO",
        },
        {
            "metro": "Memphis, TN",
            "concentration": 0.80,
            "employers": "FedEx, AutoZone DC",
        },
        {
            "metro": "Indianapolis, IN",
            "concentration": 0.78,
            "employers": "FedEx Ground, Amazon",
        },
    ],
    "Welding": [
        {
            "metro": "Houston-The Woodlands, TX",
            "concentration": 0.92,
            "employers": "Bechtel, Kiewit",
        },
        {
            "metro": "Pittsburgh, PA",
            "concentration": 0.78,
            "employers": "US Steel, BWXT",
        },
        {
            "metro": "Detroit-Warren, MI",
            "concentration": 0.75,
            "employers": "GM, Ford suppliers",
        },
        {"metro": "Tulsa, OK", "concentration": 0.72, "employers": "Hilti, ONEOK"},
    ],
    "HVAC": [
        {
            "metro": "Phoenix-Mesa, AZ",
            "concentration": 0.88,
            "employers": "Parker & Sons, Hays Cooling",
        },
        {
            "metro": "Houston-The Woodlands, TX",
            "concentration": 0.85,
            "employers": "Service Experts",
        },
        {
            "metro": "Dallas-Fort Worth, TX",
            "concentration": 0.82,
            "employers": "Aire Serv, Lennox",
        },
    ],
    "Patient Assessment": [
        {"metro": "Rochester, MN", "concentration": 0.90, "employers": "Mayo Clinic"},
        {
            "metro": "Boston-Cambridge, MA",
            "concentration": 0.88,
            "employers": "Mass General, Dana-Farber",
        },
        {
            "metro": "Houston-The Woodlands, TX",
            "concentration": 0.85,
            "employers": "MD Anderson, Memorial Hermann",
        },
        {
            "metro": "Nashville-Davidson, TN",
            "concentration": 0.82,
            "employers": "HCA, Vanderbilt",
        },
    ],
    "Financial Modeling": [
        {
            "metro": "New York-Newark, NY",
            "concentration": 0.95,
            "employers": "Goldman Sachs, JP Morgan",
        },
        {
            "metro": "Chicago-Naperville, IL",
            "concentration": 0.80,
            "employers": "Citadel, Morningstar",
        },
        {
            "metro": "Charlotte-Concord, NC",
            "concentration": 0.75,
            "employers": "Bank of America, Wells Fargo",
        },
    ],
    "Digital Marketing": [
        {
            "metro": "New York-Newark, NY",
            "concentration": 0.90,
            "employers": "Google, Dentsu, Omnicom",
        },
        {
            "metro": "San Francisco-Oakland, CA",
            "concentration": 0.85,
            "employers": "Salesforce, Adobe",
        },
        {
            "metro": "Chicago-Naperville, IL",
            "concentration": 0.75,
            "employers": "Publicis, Leo Burnett",
        },
    ],
    "Supply Chain Management": [
        {"metro": "Memphis, TN", "concentration": 0.88, "employers": "FedEx, AutoZone"},
        {
            "metro": "Chicago-Naperville, IL",
            "concentration": 0.85,
            "employers": "Caterpillar, McDonald's HQ",
        },
        {
            "metro": "Dallas-Fort Worth, TX",
            "concentration": 0.80,
            "employers": "AT&T, Texas Instruments",
        },
    ],
}


# ===========================================================================
# 5. CHANNEL RECOMMENDATION DATA
# ===========================================================================

_SKILL_CATEGORY_CHANNELS: Dict[str, Dict[str, Any]] = {
    "technology": {
        "channels": [
            {
                "name": "LinkedIn",
                "weight": 0.28,
                "reason": "Highest tech professional density",
            },
            {"name": "Dice", "weight": 0.15, "reason": "Tech-specific job board"},
            {
                "name": "GitHub Jobs / StackOverflow",
                "weight": 0.12,
                "reason": "Developer community boards",
            },
            {
                "name": "Indeed Sponsored",
                "weight": 0.18,
                "reason": "Broad reach with tech filtering",
            },
            {
                "name": "Programmatic Display",
                "weight": 0.15,
                "reason": "Retarget tech blog visitors",
            },
            {
                "name": "AngelList / Wellfound",
                "weight": 0.07,
                "reason": "Startup-oriented talent",
            },
            {
                "name": "Hired / Triplebyte",
                "weight": 0.05,
                "reason": "Pre-vetted tech talent",
            },
        ],
        "messaging": "Career growth, tech stack, remote flexibility, equity/compensation",
        "avg_cpc_range": [1.80, 5.50],
        "avg_cpa_range": [25, 80],
    },
    "healthcare": {
        "channels": [
            {
                "name": "Vivian Health",
                "weight": 0.22,
                "reason": "Nursing/allied health marketplace",
            },
            {
                "name": "Health eCareers",
                "weight": 0.18,
                "reason": "Multi-discipline healthcare board",
            },
            {
                "name": "Indeed Sponsored",
                "weight": 0.20,
                "reason": "Broad healthcare applicant pool",
            },
            {
                "name": "LinkedIn",
                "weight": 0.12,
                "reason": "Advanced practice and leadership roles",
            },
            {
                "name": "Programmatic Display",
                "weight": 0.15,
                "reason": "Target healthcare content readers",
            },
            {
                "name": "NurseFly",
                "weight": 0.08,
                "reason": "Travel nursing specialists",
            },
            {
                "name": "Facebook/Instagram",
                "weight": 0.05,
                "reason": "Community-based outreach",
            },
        ],
        "messaging": "Sign-on bonus, shift flexibility, patient ratio, CE support",
        "avg_cpc_range": [0.90, 3.50],
        "avg_cpa_range": [18, 55],
    },
    "trades": {
        "channels": [
            {
                "name": "Indeed Sponsored",
                "weight": 0.28,
                "reason": "Largest blue-collar applicant pool",
            },
            {
                "name": "Facebook Jobs",
                "weight": 0.22,
                "reason": "High mobile/trade-worker engagement",
            },
            {
                "name": "Programmatic Display",
                "weight": 0.20,
                "reason": "Geo-targeted mobile display",
            },
            {
                "name": "ZipRecruiter",
                "weight": 0.12,
                "reason": "Quick-apply for hourly roles",
            },
            {"name": "Craigslist", "weight": 0.08, "reason": "Local trades hiring"},
            {
                "name": "Talroo",
                "weight": 0.05,
                "reason": "Performance-based blue collar",
            },
            {
                "name": "Trade-specific boards",
                "weight": 0.05,
                "reason": "HVAC-Talk, WeldingWeb, etc.",
            },
        ],
        "messaging": "Pay rate, sign-on bonus, tools provided, schedule, proximity",
        "avg_cpc_range": [0.30, 1.50],
        "avg_cpa_range": [8, 28],
    },
    "business": {
        "channels": [
            {
                "name": "LinkedIn",
                "weight": 0.30,
                "reason": "Professional/business role standard",
            },
            {
                "name": "Indeed Sponsored",
                "weight": 0.22,
                "reason": "Volume + brand awareness",
            },
            {
                "name": "Glassdoor",
                "weight": 0.12,
                "reason": "Employer brand + professional search",
            },
            {
                "name": "Programmatic Display",
                "weight": 0.15,
                "reason": "Retarget business media readers",
            },
            {
                "name": "Google Search Ads",
                "weight": 0.10,
                "reason": "Intent-based capture",
            },
            {
                "name": "Niche boards",
                "weight": 0.06,
                "reason": "eFinancialCareers, MarketingHire",
            },
            {
                "name": "Employer Career Site",
                "weight": 0.05,
                "reason": "Direct applicants, lowest CPA",
            },
        ],
        "messaging": "Career progression, compensation, culture, hybrid/remote, benefits",
        "avg_cpc_range": [1.20, 4.00],
        "avg_cpa_range": [20, 65],
    },
    "service": {
        "channels": [
            {
                "name": "Indeed Sponsored",
                "weight": 0.30,
                "reason": "Volume hiring standard",
            },
            {
                "name": "Facebook Jobs",
                "weight": 0.25,
                "reason": "High mobile engagement",
            },
            {"name": "Snagajob", "weight": 0.15, "reason": "Hourly-focused platform"},
            {
                "name": "Programmatic Display",
                "weight": 0.15,
                "reason": "Geo-targeted mobile",
            },
            {
                "name": "Google Search Ads",
                "weight": 0.08,
                "reason": "Jobs near me intent",
            },
            {"name": "Craigslist", "weight": 0.07, "reason": "Local reach"},
        ],
        "messaging": "Flexible schedule, tips/hourly rate, team environment, benefits",
        "avg_cpc_range": [0.25, 1.00],
        "avg_cpa_range": [6, 22],
    },
    "education": {
        "channels": [
            {
                "name": "Indeed Sponsored",
                "weight": 0.25,
                "reason": "Broad educator reach",
            },
            {
                "name": "LinkedIn",
                "weight": 0.20,
                "reason": "Higher-ed / leadership roles",
            },
            {
                "name": "HigherEdJobs",
                "weight": 0.15,
                "reason": "University/college specific",
            },
            {"name": "SchoolSpring", "weight": 0.15, "reason": "K-12 focused board"},
            {
                "name": "Programmatic Display",
                "weight": 0.12,
                "reason": "Education media targeting",
            },
            {
                "name": "State job boards",
                "weight": 0.08,
                "reason": "Public school postings",
            },
            {"name": "Facebook", "weight": 0.05, "reason": "Community teacher groups"},
        ],
        "messaging": "Mission, impact, benefits, schedule, professional development",
        "avg_cpc_range": [0.80, 2.50],
        "avg_cpa_range": [15, 45],
    },
    "legal": {
        "channels": [
            {
                "name": "LinkedIn",
                "weight": 0.30,
                "reason": "Legal professional standard",
            },
            {"name": "LawCrossing", "weight": 0.15, "reason": "Legal niche board"},
            {
                "name": "Indeed Sponsored",
                "weight": 0.20,
                "reason": "Paralegal + associate volume",
            },
            {
                "name": "Programmatic Display",
                "weight": 0.12,
                "reason": "Legal media retargeting",
            },
            {
                "name": "Robert Half Legal",
                "weight": 0.10,
                "reason": "Legal staffing network",
            },
            {"name": "Glassdoor", "weight": 0.08, "reason": "Firm reputation research"},
            {
                "name": "Google Search Ads",
                "weight": 0.05,
                "reason": "Intent-based capture",
            },
        ],
        "messaging": "Practice area, billable expectations, partnership track, pro bono",
        "avg_cpc_range": [2.00, 6.00],
        "avg_cpa_range": [30, 90],
    },
}

_SKILL_TO_CATEGORY: Dict[str, str] = {
    "Python": "technology",
    "JavaScript": "technology",
    "SQL": "technology",
    "Cloud (AWS/GCP/Azure)": "technology",
    "Machine Learning": "technology",
    "Cybersecurity": "technology",
    "DevOps": "technology",
    "Kubernetes": "technology",
    "React": "technology",
    "TypeScript": "technology",
    "Java": "technology",
    "C++": "technology",
    "Rust": "technology",
    "Go": "technology",
    "Data Visualization": "technology",
    "TensorFlow/PyTorch": "technology",
    "System Design": "technology",
    "Git": "technology",
    "Salesforce": "technology",
    "SAP": "technology",
    "Patient Assessment": "healthcare",
    "BLS/ACLS": "healthcare",
    "Electronic Health Records": "healthcare",
    "Medication Administration": "healthcare",
    "Clinical Diagnosis": "healthcare",
    "Therapeutic Techniques": "healthcare",
    "Phlebotomy": "healthcare",
    "Medical Coding": "healthcare",
    "CDL License": "trades",
    "Welding": "trades",
    "Forklift Operation": "trades",
    "Blueprint Reading": "trades",
    "HVAC": "trades",
    "Electrical Troubleshooting": "trades",
    "CNC Machining": "trades",
    "Heavy Equipment Operation": "trades",
    "Carpentry": "trades",
    "Plumbing": "trades",
    "Project Management": "business",
    "Data Analysis": "business",
    "Financial Modeling": "business",
    "SEO/SEM": "business",
    "Digital Marketing": "business",
    "CRM Software": "business",
    "Negotiation": "business",
    "UX Design": "business",
    "Figma/Sketch": "business",
    "AutoCAD": "business",
    "Six Sigma": "business",
    "Supply Chain Management": "business",
    "Accounting (GAAP)": "business",
    "Customer Service": "service",
    "Food Safety (ServSafe)": "service",
    "POS Systems": "service",
    "Curriculum Development": "education",
    "Legal Research": "legal",
    "Contract Negotiation": "legal",
}


# ===========================================================================
# 6. INTERNAL HELPERS (defined before public functions that use them)
# ===========================================================================

_SKILL_ALIASES: Dict[str, str] = {
    "aws": "Cloud (AWS/GCP/Azure)",
    "gcp": "Cloud (AWS/GCP/Azure)",
    "azure": "Cloud (AWS/GCP/Azure)",
    "cloud": "Cloud (AWS/GCP/Azure)",
    "ml": "Machine Learning",
    "ai": "Machine Learning",
    "deep learning": "Machine Learning",
    "tensorflow": "TensorFlow/PyTorch",
    "pytorch": "TensorFlow/PyTorch",
    "torch": "TensorFlow/PyTorch",
    "k8s": "Kubernetes",
    "docker": "DevOps",
    "ci/cd": "DevOps",
    "cicd": "DevOps",
    "react.js": "React",
    "reactjs": "React",
    "node.js": "JavaScript",
    "nodejs": "JavaScript",
    "angular": "JavaScript",
    "vue": "JavaScript",
    "ts": "TypeScript",
    "figma": "Figma/Sketch",
    "sketch": "Figma/Sketch",
    "ux": "UX Design",
    "ui": "UX Design",
    "ui/ux": "UX Design",
    "gaap": "Accounting (GAAP)",
    "ifrs": "Accounting (GAAP)",
    "sap erp": "SAP",
    "oracle erp": "SAP",
    "sfdc": "Salesforce",
    "salesforce crm": "Salesforce",
    "pmp": "Project Management",
    "agile": "Project Management",
    "scrum": "Project Management",
    "seo": "SEO/SEM",
    "sem": "SEO/SEM",
    "google ads": "SEO/SEM",
    "cdl": "CDL License",
    "class a": "CDL License",
    "class a cdl": "CDL License",
    "hvac-r": "HVAC",
    "refrigeration": "HVAC",
    "mig welding": "Welding",
    "tig welding": "Welding",
    "arc welding": "Welding",
    "stick welding": "Welding",
    "cnc": "CNC Machining",
    "lathe": "CNC Machining",
    "servsafe": "Food Safety (ServSafe)",
    "autocad": "AutoCAD",
    "revit": "AutoCAD",
    "solidworks": "AutoCAD",
    "lean": "Six Sigma",
    "lean six sigma": "Six Sigma",
    "supply chain": "Supply Chain Management",
    "logistics": "Supply Chain Management",
    "infosec": "Cybersecurity",
    "information security": "Cybersecurity",
    "penetration testing": "Cybersecurity",
    "soc analyst": "Cybersecurity",
    "rn": "Patient Assessment",
    "registered nurse": "Patient Assessment",
    "ehr": "Electronic Health Records",
    "epic": "Electronic Health Records",
    "cerner": "Electronic Health Records",
    "acls": "BLS/ACLS",
    "bls": "BLS/ACLS",
    "icd-10": "Medical Coding",
    "cpt coding": "Medical Coding",
    "phlebotomist": "Phlebotomy",
}


def _normalize_skill(skill: str) -> str:
    """Normalize a skill string for lookup. Tries exact, case-insensitive, then alias."""
    s = skill.strip()
    if s in SKILLS_TO_OCCUPATIONS:
        return s
    lower = s.lower()
    for key in SKILLS_TO_OCCUPATIONS:
        if key.lower() == lower:
            return key
    alias_result = _SKILL_ALIASES.get(lower)
    if alias_result:
        return alias_result
    return s


_SALARY_ESTIMATES: Dict[str, Dict[str, int]] = {
    "software": {"p25": 90000, "median": 120000, "p75": 155000},
    "data scientist": {"p25": 95000, "median": 127000, "p75": 165000},
    "web developer": {"p25": 55000, "median": 78000, "p75": 105000},
    "network": {"p25": 65000, "median": 85000, "p75": 110000},
    "security": {"p25": 80000, "median": 107000, "p75": 140000},
    "database": {"p25": 72000, "median": 98000, "p75": 128000},
    "nurse": {"p25": 62000, "median": 81000, "p75": 101000},
    "physician": {"p25": 180000, "median": 230000, "p75": 310000},
    "pharmacist": {"p25": 110000, "median": 128000, "p75": 150000},
    "therapist": {"p25": 45000, "median": 58000, "p75": 78000},
    "truck driver": {"p25": 42000, "median": 52000, "p75": 65000},
    "welder": {"p25": 38000, "median": 47000, "p75": 60000},
    "electrician": {"p25": 45000, "median": 61000, "p75": 82000},
    "plumber": {"p25": 42000, "median": 59000, "p75": 80000},
    "hvac": {"p25": 40000, "median": 52000, "p75": 72000},
    "machinist": {"p25": 38000, "median": 48000, "p75": 62000},
    "carpenter": {"p25": 38000, "median": 52000, "p75": 70000},
    "laborer": {"p25": 30000, "median": 38000, "p75": 48000},
    "cashier": {"p25": 24000, "median": 28000, "p75": 33000},
    "customer service": {"p25": 30000, "median": 37000, "p75": 45000},
    "sales": {"p25": 40000, "median": 62000, "p75": 95000},
    "marketing": {"p25": 48000, "median": 68000, "p75": 95000},
    "accountant": {"p25": 55000, "median": 77000, "p75": 100000},
    "analyst": {"p25": 55000, "median": 72000, "p75": 95000},
    "manager": {"p25": 70000, "median": 100000, "p75": 140000},
    "engineer": {"p25": 75000, "median": 100000, "p75": 135000},
    "teacher": {"p25": 42000, "median": 58000, "p75": 75000},
    "lawyer": {"p25": 80000, "median": 127000, "p75": 190000},
    "paralegal": {"p25": 42000, "median": 56000, "p75": 75000},
    "operator": {"p25": 32000, "median": 42000, "p75": 55000},
    "cook": {"p25": 26000, "median": 32000, "p75": 40000},
    "medical records": {"p25": 36000, "median": 46000, "p75": 58000},
    "drafter": {"p25": 42000, "median": 56000, "p75": 72000},
    "industrial engineer": {"p25": 70000, "median": 92000, "p75": 118000},
    "logistician": {"p25": 55000, "median": 77000, "p75": 100000},
    "counselor": {"p25": 40000, "median": 50000, "p75": 65000},
    "graphic designer": {"p25": 40000, "median": 55000, "p75": 75000},
}


def _estimate_salary(occ_lower: str) -> Dict[str, int]:
    for keyword, data in _SALARY_ESTIMATES.items():
        if keyword in occ_lower:
            return data
    return {"p25": 40000, "median": 55000, "p75": 75000}


def _allocate_budget(budget: float, channels: List[Dict[str, Any]]) -> Dict[str, Any]:
    allocation = []
    for ch in channels:
        w = ch.get("weight") or 0
        amt = round(budget * w, 2)
        allocation.append(
            {
                "channel": ch.get("name") or "",
                "weight": w,
                "monthly_spend": amt,
                "annual_spend": round(amt * 12, 2),
            }
        )
    return {
        "monthly_budget": budget,
        "annual_budget": round(budget * 12, 2),
        "allocations": allocation,
    }


def _build_summary(skills, occupations, demand, collar_type):
    try:
        n_skills = len(skills)
        n_occs = len(occupations)
        top_occ = occupations[0]["title"] if occupations else "General"
        d_summary = demand.get("summary", {})
        shortage = d_summary.get("skills_in_shortage") or 0
        overall = d_summary.get("overall_demand", "balanced").replace("_", " ")
        collar_label = collar_type.replace("_", " ").title()
        parts = [
            f"Analyzed {n_skills} skill(s) matching {n_occs} O*NET occupation(s).",
            f"Primary occupation: {top_occ} ({collar_label}).",
            f"Market outlook: {overall} ({shortage} of {n_skills} skills in shortage).",
        ]
        if shortage > n_skills * 0.5:
            parts.append(
                "ALERT: Over half of the target skills are in shortage. "
                "Consider higher CPA budgets, sign-on bonuses, and niche channel investment."
            )
        elif shortage > 0:
            parts.append(
                "Some skills show supply constraints. "
                "Targeted sourcing on niche boards recommended for scarce skills."
            )
        else:
            parts.append(
                "Skills are broadly available. "
                "Focus on employer brand and competitive compensation."
            )
        return " ".join(parts)
    except Exception:
        return "Skill analysis complete. See detailed sections for findings."


# ===========================================================================
# 7. CORE PUBLIC FUNCTIONS
# ===========================================================================


def find_matching_occupations(skills: List[str]) -> List[Dict[str, Any]]:
    """Match skills to O*NET-aligned occupations, sorted by match count."""
    try:
        occ_map: Dict[str, Dict[str, Any]] = {}
        for skill in skills:
            skill_key = _normalize_skill(skill)
            for occ in SKILLS_TO_OCCUPATIONS.get(skill_key, []):
                soc = occ["soc"]
                if soc not in occ_map:
                    occ_map[soc] = {
                        "soc": soc,
                        "title": occ["title"],
                        "zone": occ.get("zone", "3"),
                        "matched_skills": [],
                        "match_count": 0,
                    }
                if skill_key not in occ_map[soc]["matched_skills"]:
                    occ_map[soc]["matched_skills"].append(skill_key)
                    occ_map[soc]["match_count"] += 1
        return sorted(occ_map.values(), key=lambda d: d["match_count"], reverse=True)
    except Exception as e:
        logger.error("find_matching_occupations error: %s", e)
        return []


def get_skill_demand_trends(skills: List[str]) -> Dict[str, Any]:
    """Return demand indicators for each skill."""
    try:
        results: Dict[str, Any] = {}
        shortage_count = high_growth_count = 0
        for skill in skills:
            key = _normalize_skill(skill)
            info = _SKILL_DEMAND.get(
                key,
                {
                    "growth": "unknown",
                    "yoy_pct": 0,
                    "pool": "unknown",
                    "shortage": False,
                },
            )
            results[key] = info
            if info.get("shortage"):
                shortage_count += 1
            if info.get("growth") in ("high", "very_high"):
                high_growth_count += 1
        total = max(len(skills), 1)
        return {
            "skills": results,
            "summary": {
                "total_skills_analyzed": len(skills),
                "skills_in_shortage": shortage_count,
                "shortage_pct": round(shortage_count / total * 100, 1),
                "high_growth_skills": high_growth_count,
                "high_growth_pct": round(high_growth_count / total * 100, 1),
                "overall_demand": (
                    "critical_shortage"
                    if shortage_count / total > 0.5
                    else "tight" if shortage_count / total > 0.25 else "balanced"
                ),
            },
        }
    except Exception as e:
        logger.error("get_skill_demand_trends error: %s", e)
        return {"skills": {}, "summary": {}}


def recommend_channels(skill_profile: Dict[str, Any]) -> List[Dict[str, Any]]:
    """Recommend recruiting channels based on skill profile + collar type."""
    try:
        skills = skill_profile.get("skills") or []
        collar = skill_profile.get("collar_type") or ""
        categories: Dict[str, int] = {}
        for skill in skills:
            cat = _SKILL_TO_CATEGORY.get(_normalize_skill(skill), "business")
            categories[cat] = categories.get(cat, 0) + 1
        if not categories:
            categories = {"business": 1}
        total_cat = sum(categories.values())
        cat_weights = {c: cnt / total_cat for c, cnt in categories.items()}
        channel_scores: Dict[str, Dict[str, Any]] = {}
        for cat, cat_w in cat_weights.items():
            cat_data = _SKILL_CATEGORY_CHANNELS.get(
                cat, _SKILL_CATEGORY_CHANNELS["business"]
            )
            for ch in cat_data["channels"]:
                name = ch["name"]
                if name not in channel_scores:
                    channel_scores[name] = {"name": name, "weight": 0.0, "reasons": []}
                channel_scores[name]["weight"] += ch["weight"] * cat_w
                channel_scores[name]["reasons"].append(ch["reason"])
        if collar and _HAS_COLLAR and collar in COLLAR_STRATEGY:
            for platform in (
                COLLAR_STRATEGY[collar].get("preferred_platforms") or [][:3]
            ):
                for ch_name in channel_scores:
                    if platform.lower() in ch_name.lower():
                        channel_scores[ch_name]["weight"] *= 1.15
        total_w = sum(ch["weight"] for ch in channel_scores.values())
        results = []
        for ch in channel_scores.values():
            ch["weight"] = round(ch["weight"] / total_w, 3) if total_w > 0 else 0
            ch["reasons"] = list(dict.fromkeys(ch["reasons"]))[:3]
            results.append(ch)
        results.sort(key=lambda d: d["weight"], reverse=True)
        return results[:10]
    except Exception as e:
        logger.error("recommend_channels error: %s", e)
        return []


def get_salary_benchmarks(occupations: List[str], location: str = "") -> Dict[str, Any]:
    """Get salary benchmarks, optionally COLI-adjusted by location."""
    try:
        salary_data: Dict[str, Any] = {}
        api_data = {}
        if _HAS_API and occupations:
            try:
                api_data = fetch_salary_data(occupations[:10])
            except Exception:
                pass
        coli = 100.0
        if location and _HAS_RESEARCH:
            try:
                coli = get_location_info(location).get("coli", 100.0)
            except Exception:
                pass
        coli_factor = coli / 100.0
        for occ in occupations[:15]:
            occ_lower = occ.lower().strip()
            if occ in api_data and api_data[occ].get("median"):
                raw = api_data[occ]
                salary_data[occ] = {
                    "source": "BLS API",
                    "median": round(raw["median"] * coli_factor),
                    "p25": round(raw.get("p25", raw["median"] * 0.78) * coli_factor),
                    "p75": round(raw.get("p75", raw["median"] * 1.25) * coli_factor),
                    "coli_adjusted": coli != 100.0,
                    "coli": round(coli, 1),
                }
            else:
                est = _estimate_salary(occ_lower)
                salary_data[occ] = {
                    "source": "built-in estimate",
                    "median": round(est["median"] * coli_factor),
                    "p25": round(est["p25"] * coli_factor),
                    "p75": round(est["p75"] * coli_factor),
                    "coli_adjusted": coli != 100.0,
                    "coli": round(coli, 1),
                }
        medians = [v["median"] for v in salary_data.values() if v.get("median")]
        return {
            "occupations": salary_data,
            "aggregate": {
                "avg_median": round(sum(medians) / len(medians)) if medians else 0,
                "min_median": min(medians) if medians else 0,
                "max_median": max(medians) if medians else 0,
                "location": location or "National (no location specified)",
                "coli": round(coli, 1),
            },
        }
    except Exception as e:
        logger.error("get_salary_benchmarks error: %s", e)
        return {"occupations": {}, "aggregate": {}}


def find_skill_hotspots(skills: List[str]) -> List[Dict[str, Any]]:
    """Find geographic concentrations for skills."""
    try:
        metro_scores: Dict[str, Dict[str, Any]] = {}
        for skill in skills:
            key = _normalize_skill(skill)
            for hs in _SKILL_HOTSPOTS.get(key, []):
                metro = hs["metro"]
                if metro not in metro_scores:
                    metro_scores[metro] = {
                        "metro": metro,
                        "concentrations": [],
                        "skills_present": [],
                        "top_employers": set(),
                    }
                metro_scores[metro]["concentrations"].append(hs["concentration"])
                metro_scores[metro]["skills_present"].append(key)
                for emp in hs.get("employers") or "".split(", "):
                    if emp.strip():
                        metro_scores[metro]["top_employers"].add(emp.strip())
        if _HAS_RESEARCH and not metro_scores:
            for metro in [
                "New York, NY",
                "Los Angeles, CA",
                "Chicago, IL",
                "Houston, TX",
                "Dallas, TX",
            ]:
                try:
                    info = get_location_info(metro)
                    metro_scores[metro] = {
                        "metro": metro,
                        "concentrations": [0.50],
                        "skills_present": list(skills[:3]),
                        "top_employers": set(
                            info.get("major_employers") or "".split(", ")[:3]
                        ),
                    }
                except Exception:
                    pass
        results = []
        for metro, data in metro_scores.items():
            concs = data["concentrations"]
            results.append(
                {
                    "metro": data["metro"],
                    "avg_concentration": (
                        round(sum(concs) / len(concs), 3) if concs else 0
                    ),
                    "skills_present": list(dict.fromkeys(data["skills_present"])),
                    "skill_coverage": round(
                        len(set(data["skills_present"])) / max(len(skills), 1), 2
                    ),
                    "top_employers": sorted(data["top_employers"])[:6],
                }
            )
        results.sort(key=lambda d: d["avg_concentration"], reverse=True)
        return results[:15]
    except Exception as e:
        logger.error("find_skill_hotspots error: %s", e)
        return []


def get_adjacent_skills(skills: List[str]) -> List[Dict[str, Any]]:
    """Find related skills, excluding those already in the input set."""
    try:
        input_set = {_normalize_skill(s) for s in skills}
        adj_map: Dict[str, Dict[str, Any]] = {}
        for skill in skills:
            key = _normalize_skill(skill)
            for adj in ADJACENT_SKILLS.get(key, []):
                adj_name = adj["skill"]
                if adj_name in input_set:
                    continue
                if adj_name not in adj_map:
                    adj_map[adj_name] = {
                        "skill": adj_name,
                        "max_relevance": 0.0,
                        "connected_to": [],
                        "scarcity": (
                            SKILL_SCARCITY.get(adj_name, 0.40) if _HAS_COLLAR else 0.40
                        ),
                    }
                if adj["relevance"] > adj_map[adj_name]["max_relevance"]:
                    adj_map[adj_name]["max_relevance"] = adj["relevance"]
                adj_map[adj_name]["connected_to"].append(key)
        results = sorted(
            adj_map.values(), key=lambda d: d["max_relevance"], reverse=True
        )
        for r in results:
            r["max_relevance"] = round(r["max_relevance"], 2)
            r["connected_to"] = list(dict.fromkeys(r["connected_to"]))
        return results[:20]
    except Exception as e:
        logger.error("get_adjacent_skills error: %s", e)
        return []


def suggest_job_titles(skills: List[str], industry: str = "") -> List[str]:
    """Generate optimized job title suggestions."""
    try:
        occupations = find_matching_occupations(skills)
        base_titles = [occ["title"] for occ in occupations[:8]]
        enhanced: List[str] = []
        for title in base_titles:
            enhanced.append(title)
            if industry:
                label = INDUSTRY_LABEL_MAP.get(industry, "") if _HAS_UTILS else ""
                if not label:
                    label = industry.replace("_", " ").title()
                if len(label) < 25:
                    enhanced.append(f"{title} ({label})")
        if _HAS_COLLAR:
            categories = {}
            for s in skills:
                cat = _SKILL_TO_CATEGORY.get(_normalize_skill(s), "business")
                categories[cat] = categories.get(cat, 0) + 1
            primary_cat = (
                max(categories, key=categories.get) if categories else "business"
            )
            suffixes = {
                "technology": ["(Remote Available)", "(Hybrid)"],
                "healthcare": ["(Day Shift)", "(Night Shift)", "(PRN)"],
                "trades": ["(Full-Time)", "(Overtime Available)"],
                "business": ["(Hybrid)", "(On-Site)"],
                "service": ["(Full-Time)", "(Part-Time)"],
            }
            for suffix in suffixes.get(primary_cat, []):
                if base_titles:
                    enhanced.append(f"{base_titles[0]} {suffix}")
        seen = set()
        unique = []
        for t in enhanced:
            if t.lower() not in seen:
                seen.add(t.lower())
                unique.append(t)
        return unique[:15]
    except Exception as e:
        logger.error("suggest_job_titles error: %s", e)
        return []


def analyze_skills(
    skills: List[str], industry: str = "", location: str = ""
) -> Dict[str, Any]:
    """Main orchestrator: analyze skills for recruitment targeting."""
    try:
        occupations = find_matching_occupations(skills)
        demand = get_skill_demand_trends(skills)
        adjacent = get_adjacent_skills(skills)
        hotspots = find_skill_hotspots(skills)
        titles = suggest_job_titles(skills, industry)
        collar_type = "white_collar"
        if _HAS_COLLAR and occupations:
            collar_type = classify_collar(occupations[0]["title"], industry).get(
                "collar_type", "white_collar"
            )
        channels = recommend_channels({"skills": skills, "collar_type": collar_type})
        salary = get_salary_benchmarks([o["title"] for o in occupations[:5]], location)
        return {
            "input": {"skills": skills, "industry": industry, "location": location},
            "occupations": occupations,
            "demand_trends": demand,
            "channels": channels,
            "salary_benchmarks": salary,
            "hotspots": hotspots,
            "adjacent_skills": adjacent,
            "suggested_titles": titles,
            "collar_type": collar_type,
            "summary": _build_summary(skills, occupations, demand, collar_type),
        }
    except Exception as e:
        logger.error("analyze_skills error: %s", e)
        return {"error": str(e), "input": {"skills": skills}}


# ===========================================================================
# 8. FULL ANALYSIS ORCHESTRATOR (threaded)
# ===========================================================================


def run_full_skill_analysis(
    skills: List[str],
    industry: str = "",
    location: str = "",
    budget: Optional[float] = None,
) -> Dict[str, Any]:
    """Run comprehensive skill analysis with parallel sub-tasks."""
    try:
        if not skills:
            return {"error": "No skills provided", "input": {"skills": []}}
        skills = [s.strip() for s in skills if s.strip()]
        occupations = find_matching_occupations(skills)
        collar_type = "white_collar"
        if _HAS_COLLAR and occupations:
            collar_type = classify_collar(occupations[0]["title"], industry).get(
                "collar_type", "white_collar"
            )
        skill_profile = {"skills": skills, "collar_type": collar_type}
        results: Dict[str, Any] = {
            "occupations": occupations,
            "collar_type": collar_type,
        }
        with ThreadPoolExecutor(max_workers=5) as executor:
            futures = {
                executor.submit(get_skill_demand_trends, skills): "demand_trends",
                executor.submit(recommend_channels, skill_profile): "channels",
                executor.submit(
                    get_salary_benchmarks,
                    [o["title"] for o in occupations[:5]],
                    location,
                ): "salary_benchmarks",
                executor.submit(find_skill_hotspots, skills): "hotspots",
                executor.submit(get_adjacent_skills, skills): "adjacent_skills",
            }
            for future in as_completed(futures):
                key = futures[future]
                try:
                    results[key] = future.result(timeout=15)
                except Exception as exc:
                    logger.warning("Sub-task %s failed: %s", key, exc)
                    results[key] = {} if key.endswith("trends") else []
        results["suggested_titles"] = suggest_job_titles(skills, industry)
        if budget and budget > 0:
            results["budget_allocation"] = _allocate_budget(
                budget, results.get("channels") or []
            )
        if _HAS_COLLAR and occupations:
            try:
                results["skills_gap"] = analyze_skills_gap(
                    occupations[0]["title"], location, industry
                )
            except Exception:
                pass
        results["input"] = {
            "skills": skills,
            "industry": industry,
            "location": location,
            "budget": budget,
        }
        results["summary"] = _build_summary(
            skills, occupations, results.get("demand_trends", {}), collar_type
        )
        return results
    except Exception as e:
        logger.error("run_full_skill_analysis error: %s", e)
        return {"error": str(e), "input": {"skills": skills}}


# ===========================================================================
# 9. EXCEL EXPORT  (Sapphire Blue palette)
# ===========================================================================


def generate_skill_excel(analysis: Dict[str, Any]) -> bytes:
    """Generate formatted Excel. Palette: #0F172A, #2563EB, #DBEAFE. Calibri, col B."""
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    except ImportError:
        logger.error("openpyxl not installed")
        return b""
    try:
        wb = Workbook()
        DARK, ACCENT, LIGHT, WHITE = "0F172A", "2563EB", "DBEAFE", "FFFFFF"
        hdr_font = Font(name="Calibri", bold=True, color=WHITE, size=12)
        body_font = Font(name="Calibri", color=DARK, size=10)
        accent_font = Font(name="Calibri", bold=True, color=ACCENT, size=11)
        dark_fill = PatternFill(start_color=DARK, end_color=DARK, fill_type="solid")
        light_fill = PatternFill(start_color=LIGHT, end_color=LIGHT, fill_type="solid")
        bdr = Border(
            left=Side(style="thin", color=ACCENT),
            right=Side(style="thin", color=ACCENT),
            top=Side(style="thin", color=ACCENT),
            bottom=Side(style="thin", color=ACCENT),
        )

        def _hdr(ws, row, cols):
            for c in range(2, 2 + cols):
                cell = ws.cell(row=row, column=c)
                cell.fill = dark_fill
                cell.font = hdr_font
                cell.alignment = Alignment(horizontal="center", vertical="center")
                cell.border = bdr

        def _drow(ws, row, cols, alt=False):
            for c in range(2, 2 + cols):
                cell = ws.cell(row=row, column=c)
                cell.font = body_font
                cell.border = bdr
                cell.alignment = Alignment(vertical="center", wrap_text=True)
                if alt:
                    cell.fill = light_fill

        # Sheet 1: Occupations
        ws1 = wb.active
        ws1.title = "Matched Occupations"
        ws1.sheet_properties.tabColor = ACCENT
        ws1.column_dimensions["A"].width = 3
        ws1.column_dimensions["B"].width = 14
        ws1.column_dimensions["C"].width = 38
        ws1.column_dimensions["D"].width = 12
        ws1.column_dimensions["E"].width = 14
        ws1.column_dimensions["F"].width = 45
        ws1.merge_cells("B1:F1")
        ws1["B1"].value = "Skill-Based Occupation Matching"
        ws1["B1"].font = Font(name="Calibri", bold=True, color=WHITE, size=14)
        ws1["B1"].fill = dark_fill
        ws1["B1"].alignment = Alignment(horizontal="center")
        inp = analysis.get("input", {})
        ws1["B2"] = f"Skills: {', '.join(inp.get('skills') or [])}"
        ws1["B2"].font = accent_font
        ws1["B3"] = (
            f"Industry: {inp.get('industry', 'N/A')}  |  Location: {inp.get('location', 'N/A')}"
        )
        ws1["B3"].font = body_font
        row = 5
        hdrs = [
            "SOC Code",
            "Occupation Title",
            "Job Zone",
            "Match Count",
            "Matched Skills",
        ]
        for i, h in enumerate(hdrs):
            ws1.cell(row=row, column=2 + i, value=h)
        _hdr(ws1, row, len(hdrs))
        for idx, occ in enumerate(analysis.get("occupations") or [][:20]):
            r = row + 1 + idx
            ws1.cell(row=r, column=2, value=occ.get("soc") or "")
            ws1.cell(row=r, column=3, value=occ.get("title") or "")
            ws1.cell(row=r, column=4, value=occ.get("zone") or "")
            ws1.cell(row=r, column=5, value=occ.get("match_count") or 0)
            ws1.cell(row=r, column=6, value=", ".join(occ.get("matched_skills") or []))
            _drow(ws1, r, len(hdrs), alt=(idx % 2 == 0))

        # Sheet 2: Demand
        ws2 = wb.create_sheet("Demand Trends")
        ws2.sheet_properties.tabColor = ACCENT
        for col, w in [("A", 3), ("B", 32), ("C", 14), ("D", 14), ("E", 14), ("F", 14)]:
            ws2.column_dimensions[col].width = w
        ws2.merge_cells("B1:F1")
        ws2["B1"].value = "Skill Demand Trends"
        ws2["B1"].font = Font(name="Calibri", bold=True, color=WHITE, size=14)
        ws2["B1"].fill = dark_fill
        ws2["B1"].alignment = Alignment(horizontal="center")
        row = 3
        dh = ["Skill", "Growth", "YoY %", "Pool Size", "Shortage?"]
        for i, h in enumerate(dh):
            ws2.cell(row=row, column=2 + i, value=h)
        _hdr(ws2, row, len(dh))
        for idx, (sk, info) in enumerate(
            analysis.get("demand_trends", {}).get("skills", {}).items()
        ):
            r = row + 1 + idx
            ws2.cell(row=r, column=2, value=sk)
            ws2.cell(row=r, column=3, value=info.get("growth") or "")
            ws2.cell(row=r, column=4, value=f"{info.get('yoy_pct') or 0}%")
            ws2.cell(row=r, column=5, value=info.get("pool") or "")
            ws2.cell(row=r, column=6, value="YES" if info.get("shortage") else "No")
            _drow(ws2, r, len(dh), alt=(idx % 2 == 0))

        # Sheet 3: Channels
        ws3 = wb.create_sheet("Channels")
        ws3.sheet_properties.tabColor = ACCENT
        for col, w in [("A", 3), ("B", 32), ("C", 14), ("D", 50)]:
            ws3.column_dimensions[col].width = w
        ws3.merge_cells("B1:D1")
        ws3["B1"].value = "Recommended Channels"
        ws3["B1"].font = Font(name="Calibri", bold=True, color=WHITE, size=14)
        ws3["B1"].fill = dark_fill
        ws3["B1"].alignment = Alignment(horizontal="center")
        row = 3
        ch_h = ["Channel", "Weight", "Rationale"]
        for i, h in enumerate(ch_h):
            ws3.cell(row=row, column=2 + i, value=h)
        _hdr(ws3, row, len(ch_h))
        for idx, ch in enumerate(analysis.get("channels") or []):
            r = row + 1 + idx
            ws3.cell(row=r, column=2, value=ch.get("name") or "")
            ws3.cell(row=r, column=3, value=f"{ch.get('weight') or 0:.1%}")
            ws3.cell(row=r, column=4, value="; ".join(ch.get("reasons") or []))
            _drow(ws3, r, len(ch_h), alt=(idx % 2 == 0))

        # Sheet 4: Salary
        ws4 = wb.create_sheet("Salary Benchmarks")
        ws4.sheet_properties.tabColor = ACCENT
        for col, w in [("A", 3), ("B", 38), ("C", 16), ("D", 16), ("E", 16), ("F", 14)]:
            ws4.column_dimensions[col].width = w
        ws4.merge_cells("B1:F1")
        ws4["B1"].value = "Salary Benchmarks"
        ws4["B1"].font = Font(name="Calibri", bold=True, color=WHITE, size=14)
        ws4["B1"].fill = dark_fill
        ws4["B1"].alignment = Alignment(horizontal="center")
        row = 3
        sh = ["Occupation", "25th %ile", "Median", "75th %ile", "Source"]
        for i, h in enumerate(sh):
            ws4.cell(row=row, column=2 + i, value=h)
        _hdr(ws4, row, len(sh))
        for idx, (occ, info) in enumerate(
            analysis.get("salary_benchmarks", {}).get("occupations", {}).items()
        ):
            r = row + 1 + idx
            ws4.cell(row=r, column=2, value=occ)
            ws4.cell(row=r, column=3, value=f"${info.get('p25') or 0:,}")
            ws4.cell(row=r, column=4, value=f"${info.get('median') or 0:,}")
            ws4.cell(row=r, column=5, value=f"${info.get('p75') or 0:,}")
            ws4.cell(row=r, column=6, value=info.get("source") or "")
            _drow(ws4, r, len(sh), alt=(idx % 2 == 0))

        # Sheet 5: Hotspots
        ws5 = wb.create_sheet("Hotspots")
        ws5.sheet_properties.tabColor = ACCENT
        for col, w in [("A", 3), ("B", 36), ("C", 18), ("D", 16), ("E", 42)]:
            ws5.column_dimensions[col].width = w
        ws5.merge_cells("B1:E1")
        ws5["B1"].value = "Skill Geographic Hotspots"
        ws5["B1"].font = Font(name="Calibri", bold=True, color=WHITE, size=14)
        ws5["B1"].fill = dark_fill
        ws5["B1"].alignment = Alignment(horizontal="center")
        row = 3
        hh = ["Metro Area", "Concentration", "Skill Coverage", "Top Employers"]
        for i, h in enumerate(hh):
            ws5.cell(row=row, column=2 + i, value=h)
        _hdr(ws5, row, len(hh))
        for idx, hs in enumerate(analysis.get("hotspots") or [][:15]):
            r = row + 1 + idx
            ws5.cell(row=r, column=2, value=hs.get("metro") or "")
            ws5.cell(row=r, column=3, value=f"{hs.get('avg_concentration') or 0:.0%}")
            ws5.cell(row=r, column=4, value=f"{hs.get('skill_coverage') or 0:.0%}")
            ws5.cell(row=r, column=5, value=", ".join(hs.get("top_employers") or []))
            _drow(ws5, r, len(hh), alt=(idx % 2 == 0))

        # Sheet 6: Adjacent Skills
        ws6 = wb.create_sheet("Adjacent Skills")
        ws6.sheet_properties.tabColor = ACCENT
        for col, w in [("A", 3), ("B", 30), ("C", 14), ("D", 14), ("E", 40)]:
            ws6.column_dimensions[col].width = w
        ws6.merge_cells("B1:E1")
        ws6["B1"].value = "Adjacent / Related Skills"
        ws6["B1"].font = Font(name="Calibri", bold=True, color=WHITE, size=14)
        ws6["B1"].fill = dark_fill
        ws6["B1"].alignment = Alignment(horizontal="center")
        row = 3
        ah = ["Skill", "Relevance", "Scarcity", "Connected To"]
        for i, h in enumerate(ah):
            ws6.cell(row=row, column=2 + i, value=h)
        _hdr(ws6, row, len(ah))
        for idx, adj in enumerate(analysis.get("adjacent_skills") or [][:20]):
            r = row + 1 + idx
            ws6.cell(row=r, column=2, value=adj.get("skill") or "")
            ws6.cell(row=r, column=3, value=f"{adj.get('max_relevance') or 0:.0%}")
            ws6.cell(row=r, column=4, value=f"{adj.get('scarcity') or 0:.0%}")
            ws6.cell(row=r, column=5, value=", ".join(adj.get("connected_to") or []))
            _drow(ws6, r, len(ah), alt=(idx % 2 == 0))

        buf = io.BytesIO()
        wb.save(buf)
        return buf.getvalue()
    except Exception as e:
        logger.error("generate_skill_excel error: %s", e)
        return b""


# ===========================================================================
# 10. PPT EXPORT  (Nova AI Suite branding)
# ===========================================================================


def generate_skill_ppt(analysis: Dict[str, Any]) -> bytes:
    """Generate branded PPT. Port Gore #202058, Blue Violet #5A54BD, Downy #6BB3CD."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        logger.error("python-pptx not installed")
        return b""
    try:
        PG = RGBColor(0x20, 0x20, 0x58)
        BV = RGBColor(0x5A, 0x54, 0xBD)
        DT = RGBColor(0x6B, 0xB3, 0xCD)
        WH = RGBColor(0xFF, 0xFF, 0xFF)
        LG = RGBColor(0xF1, 0xF5, 0xF9)
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        def _bg(slide, color=PG):
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = color

        def _tb(slide, l, t, w, h, txt, sz=14, col=WH, bold=False, align=PP_ALIGN.LEFT):
            tx = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
            tx.text_frame.word_wrap = True
            p = tx.text_frame.paragraphs[0]
            p.text = txt
            p.font.size = Pt(sz)
            p.font.color.rgb = col
            p.font.bold = bold
            p.alignment = align
            return tx

        def _tbl(slide, l, t, w, h, rows, cols):
            return slide.shapes.add_table(
                rows, cols, Inches(l), Inches(t), Inches(w), Inches(h)
            ).table

        def _thdr(table, headers, bg=BV):
            for i, h in enumerate(headers):
                c = table.cell(0, i)
                c.text = h
                for p in c.text_frame.paragraphs:
                    p.font.size = Pt(11)
                    p.font.bold = True
                    p.font.color.rgb = WH
                    p.alignment = PP_ALIGN.CENTER
                c.fill.solid()
                c.fill.fore_color.rgb = bg

        def _tc(cell, text, alt=False):
            cell.text = str(text)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = PG
            if alt:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LG

        inp = analysis.get("input", {})
        skills = inp.get("skills") or []
        summary = analysis.get("summary") or ""

        # Slide 1: Title
        s1 = prs.slides.add_slide(prs.slide_layouts[6])
        _bg(s1)
        _tb(
            s1,
            1,
            1.5,
            11,
            1,
            "Skill-Based Targeting Analysis",
            36,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        sk_txt = ", ".join(skills[:8])
        if len(skills) > 8:
            sk_txt += f" (+{len(skills) - 8} more)"
        _tb(s1, 1, 3.0, 11, 0.8, f"Skills: {sk_txt}", 18, DT, align=PP_ALIGN.CENTER)
        meta = []
        if inp.get("industry"):
            lbl = (
                INDUSTRY_LABEL_MAP.get(inp["industry"], inp["industry"])
                if _HAS_UTILS
                else inp["industry"]
            )
            meta.append(f"Industry: {lbl}")
        if inp.get("location"):
            meta.append(f"Location: {inp['location']}")
        if inp.get("budget"):
            meta.append(f"Budget: ${inp['budget']:,.0f}/mo")
        if meta:
            _tb(s1, 1, 4.0, 11, 0.6, "  |  ".join(meta), 14, align=PP_ALIGN.CENTER)
        _tb(
            s1,
            2,
            5.5,
            9,
            0.5,
            "Powered by Nova AI Suite - SkillTarget",
            11,
            DT,
            align=PP_ALIGN.CENTER,
        )

        # Slide 2: Occupations
        s2 = prs.slides.add_slide(prs.slide_layouts[6])
        _bg(s2, WH)
        _tb(s2, 0.5, 0.3, 12, 0.6, "Matching Occupations & Summary", 24, PG, True)
        if summary:
            _tb(s2, 0.5, 1.0, 12, 0.8, summary, 12, BV)
        occs = analysis.get("occupations") or [][:8]
        if occs:
            t = _tbl(
                s2, 0.5, 2.0, 12, min((len(occs) + 1) * 0.4, 4.5), len(occs) + 1, 4
            )
            _thdr(t, ["SOC", "Occupation", "Zone", "Matched Skills"])
            for ci, w in enumerate(
                [Inches(1.5), Inches(4.5), Inches(1.2), Inches(4.8)]
            ):
                t.columns[ci].width = w
            for i, o in enumerate(occs):
                _tc(t.cell(i + 1, 0), o.get("soc") or "", i % 2 == 0)
                _tc(t.cell(i + 1, 1), o.get("title") or "", i % 2 == 0)
                _tc(t.cell(i + 1, 2), o.get("zone") or "", i % 2 == 0)
                _tc(
                    t.cell(i + 1, 3),
                    ", ".join(o.get("matched_skills") or []),
                    i % 2 == 0,
                )

        # Slide 3: Demand
        s3 = prs.slides.add_slide(prs.slide_layouts[6])
        _bg(s3, WH)
        _tb(s3, 0.5, 0.3, 12, 0.6, "Skill Demand Trends", 24, PG, True)
        ds = analysis.get("demand_trends", {}).get("summary", {})
        if ds:
            _tb(
                s3,
                0.5,
                1.0,
                12,
                0.5,
                f"Overall: {ds.get('overall_demand', 'N/A').replace('_', ' ').title()}  |  "
                f"{ds.get('skills_in_shortage') or 0} in shortage  |  "
                f"{ds.get('high_growth_skills') or 0} high-growth",
                12,
                BV,
            )
        items = list(analysis.get("demand_trends", {}).get("skills", {}).items())[:10]
        if items:
            t = _tbl(
                s3, 0.5, 1.8, 12, min((len(items) + 1) * 0.4, 4.5), len(items) + 1, 4
            )
            _thdr(t, ["Skill", "Growth", "YoY Change", "Shortage"])
            for ci, w in enumerate([Inches(4), Inches(3), Inches(2.5), Inches(2.5)]):
                t.columns[ci].width = w
            for i, (sk, info) in enumerate(items):
                _tc(t.cell(i + 1, 0), sk, i % 2 == 0)
                _tc(
                    t.cell(i + 1, 1),
                    info.get("growth") or "".replace("_", " ").title(),
                    i % 2 == 0,
                )
                _tc(t.cell(i + 1, 2), f"+{info.get('yoy_pct') or 0}%", i % 2 == 0)
                _tc(
                    t.cell(i + 1, 3),
                    "YES" if info.get("shortage") else "No",
                    i % 2 == 0,
                )

        # Slide 4: Channels
        s4 = prs.slides.add_slide(prs.slide_layouts[6])
        _bg(s4, WH)
        _tb(s4, 0.5, 0.3, 12, 0.6, "Recommended Recruiting Channels", 24, PG, True)
        chs = analysis.get("channels") or [][:8]
        if chs:
            t = _tbl(s4, 0.5, 1.2, 12, min((len(chs) + 1) * 0.45, 5), len(chs) + 1, 3)
            _thdr(t, ["Channel", "Weight", "Rationale"])
            for ci, w in enumerate([Inches(3.5), Inches(1.5), Inches(7)]):
                t.columns[ci].width = w
            for i, ch in enumerate(chs):
                _tc(t.cell(i + 1, 0), ch.get("name") or "", i % 2 == 0)
                _tc(t.cell(i + 1, 1), f"{ch.get('weight') or 0:.1%}", i % 2 == 0)
                _tc(t.cell(i + 1, 2), "; ".join(ch.get("reasons") or []), i % 2 == 0)

        # Slide 5: Salary + Hotspots
        s5 = prs.slides.add_slide(prs.slide_layouts[6])
        _bg(s5, WH)
        _tb(
            s5,
            0.5,
            0.3,
            12,
            0.6,
            "Salary Benchmarks & Geographic Hotspots",
            24,
            PG,
            True,
        )
        si = list(analysis.get("salary_benchmarks", {}).get("occupations", {}).items())[
            :5
        ]
        if si:
            t = _tbl(s5, 0.5, 1.0, 6, min((len(si) + 1) * 0.4, 2.5), len(si) + 1, 4)
            _thdr(t, ["Occupation", "P25", "Median", "P75"])
            for ci, w in enumerate(
                [Inches(2.5), Inches(1.1), Inches(1.2), Inches(1.2)]
            ):
                t.columns[ci].width = w
            for i, (occ, info) in enumerate(si):
                _tc(t.cell(i + 1, 0), occ, i % 2 == 0)
                _tc(t.cell(i + 1, 1), f"${info.get('p25') or 0:,}", i % 2 == 0)
                _tc(t.cell(i + 1, 2), f"${info.get('median') or 0:,}", i % 2 == 0)
                _tc(t.cell(i + 1, 3), f"${info.get('p75') or 0:,}", i % 2 == 0)
        hss = analysis.get("hotspots") or [][:5]
        if hss:
            ht = 1.0 + (len(si) + 1) * 0.4 + 0.5 if si else 1.5
            t = _tbl(
                s5,
                0.5,
                min(ht, 4.0),
                12,
                min((len(hss) + 1) * 0.4, 2.5),
                len(hss) + 1,
                4,
            )
            _thdr(t, ["Metro", "Concentration", "Coverage", "Top Employers"], DT)
            for ci, w in enumerate([Inches(4), Inches(2), Inches(2), Inches(4)]):
                t.columns[ci].width = w
            for i, hs in enumerate(hss):
                _tc(t.cell(i + 1, 0), hs.get("metro") or "", i % 2 == 0)
                _tc(
                    t.cell(i + 1, 1),
                    f"{hs.get('avg_concentration') or 0:.0%}",
                    i % 2 == 0,
                )
                _tc(
                    t.cell(i + 1, 2), f"{hs.get('skill_coverage') or 0:.0%}", i % 2 == 0
                )
                _tc(
                    t.cell(i + 1, 3),
                    ", ".join(hs.get("top_employers") or [][:4]),
                    i % 2 == 0,
                )

        # Slide 6: Adjacent + Titles
        s6 = prs.slides.add_slide(prs.slide_layouts[6])
        _bg(s6, WH)
        _tb(
            s6,
            0.5,
            0.3,
            12,
            0.6,
            "Adjacent Skills & Suggested Job Titles",
            24,
            PG,
            True,
        )
        adj = analysis.get("adjacent_skills") or [][:8]
        if adj:
            t = _tbl(s6, 0.5, 1.0, 6.5, min((len(adj) + 1) * 0.4, 4), len(adj) + 1, 3)
            _thdr(t, ["Related Skill", "Relevance", "Connected To"])
            for ci, w in enumerate([Inches(2.2), Inches(1.3), Inches(3.0)]):
                t.columns[ci].width = w
            for i, a in enumerate(adj):
                _tc(t.cell(i + 1, 0), a.get("skill") or "", i % 2 == 0)
                _tc(t.cell(i + 1, 1), f"{a.get('max_relevance') or 0:.0%}", i % 2 == 0)
                _tc(
                    t.cell(i + 1, 2), ", ".join(a.get("connected_to") or []), i % 2 == 0
                )
        titles = analysis.get("suggested_titles") or [][:10]
        if titles:
            _tb(s6, 7.5, 1.0, 5, 0.5, "Suggested Job Titles", 16, PG, True)
            _tb(
                s6,
                7.5,
                1.6,
                5,
                4.5,
                "\n".join(f"  {i + 1}. {t}" for i, t in enumerate(titles)),
                11,
                PG,
            )

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()
    except Exception as e:
        logger.error("generate_skill_ppt error: %s", e)
        return b""
