#!/usr/bin/env python3
"""
hire_signal.py -- HireSignal: Quality-of-Hire Signal Tracker

Extends the performance_tracker pattern to track "Quality of Hire by Source" --
the #1 metric TA leaders want. Helps recruiters understand which channels
produce the best hires, not just the most applications.

Capabilities:
  - Quality of Hire (QoH) composite score per source/channel
  - Source effectiveness: applications -> interviews -> offers -> hires -> retained funnel
  - Time-to-productivity by source
  - Cost-per-quality-hire (CPQH) vs traditional cost-per-hire
  - 90-day / 180-day / 1-year retention rates by source
  - Hiring manager satisfaction correlation
  - Diversity metrics by source
  - Predicted future performance based on source patterns
  - Excel & PowerPoint branded report generation

Thread-safe, never crashes (all errors return structured error dicts).

Design tokens:
    Excel: Sapphire Blue palette (Navy #0F172A, Sapphire #2563EB, Light #DBEAFE)
    PPT:   Joveo brand (Port Gore #202058, Blue Violet #5A54BD, Downy #6BB3CD)
"""

from __future__ import annotations

import io
import logging
import re
import random
import datetime
from concurrent.futures import ThreadPoolExecutor, as_completed
from typing import Any, Dict, List, Optional

logger = logging.getLogger(__name__)

# ── Optional imports (lazy, with try/except like performance_tracker.py) ─────
try:
    import trend_engine as _trend_engine

    _HAS_TREND_ENGINE = True
except ImportError:
    _trend_engine = None
    _HAS_TREND_ENGINE = False

try:
    from shared_utils import INDUSTRY_LABEL_MAP
except ImportError:
    INDUSTRY_LABEL_MAP = {}

    def parse_budget(v, *, default=100_000.0):
        try:
            return float(v)
        except Exception:
            return default


try:
    import performance_tracker as _perf_tracker

    _HAS_PERF_TRACKER = True
except ImportError:
    _perf_tracker = None
    _HAS_PERF_TRACKER = False


# =============================================================================
# CONSTANTS & DESIGN TOKENS
# =============================================================================

# Excel palette - Sapphire Blue
_NAVY = "0F172A"
_SAPPHIRE = "2563EB"
_BLUE_LIGHT = "DBEAFE"
_BLUE_PALE = "EFF6FF"
_WHITE = "FFFFFF"
_GREEN = "16A34A"
_GREEN_BG = "DCFCE7"
_AMBER = "D97706"
_AMBER_BG = "FEF3C7"
_RED = "DC2626"
_RED_BG = "FEE2E2"
_MUTED = "78716C"
_WARM_GRAY = "E7E5E4"
_STONE = "1C1917"

COL_START = 2  # Data starts at column B

# QoH score weights (must sum to 1.0)
_QOH_WEIGHTS = {
    "performance_rating": 0.30,
    "retention_90d": 0.20,
    "retention_180d": 0.10,
    "time_to_productivity": 0.15,
    "hiring_manager_satisfaction": 0.15,
    "cultural_fit": 0.10,
}

# Industry benchmark retention rates (fallback)
_INDUSTRY_RETENTION_BENCHMARKS = {
    "tech_engineering": {"90d": 0.88, "180d": 0.82, "1y": 0.75},
    "healthcare_medical": {"90d": 0.85, "180d": 0.78, "1y": 0.70},
    "finance_banking": {"90d": 0.90, "180d": 0.85, "1y": 0.78},
    "retail_consumer": {"90d": 0.72, "180d": 0.60, "1y": 0.48},
    "logistics_supply_chain": {"90d": 0.70, "180d": 0.58, "1y": 0.45},
    "hospitality_travel": {"90d": 0.68, "180d": 0.55, "1y": 0.42},
    "general_entry_level": {"90d": 0.75, "180d": 0.65, "1y": 0.55},
    "blue_collar_trades": {"90d": 0.73, "180d": 0.62, "1y": 0.50},
    "aerospace_defense": {"90d": 0.92, "180d": 0.88, "1y": 0.82},
    "construction_real_estate": {"90d": 0.74, "180d": 0.63, "1y": 0.52},
    "automotive": {"90d": 0.78, "180d": 0.70, "1y": 0.60},
    "energy_utilities": {"90d": 0.88, "180d": 0.82, "1y": 0.76},
    "pharma_biotech": {"90d": 0.90, "180d": 0.85, "1y": 0.78},
    "insurance": {"90d": 0.85, "180d": 0.78, "1y": 0.70},
    "telecommunications": {"90d": 0.82, "180d": 0.74, "1y": 0.65},
    "food_beverage": {"90d": 0.70, "180d": 0.57, "1y": 0.44},
    "education": {"90d": 0.88, "180d": 0.82, "1y": 0.76},
    "media_entertainment": {"90d": 0.80, "180d": 0.72, "1y": 0.62},
    "legal_services": {"90d": 0.88, "180d": 0.82, "1y": 0.75},
    "mental_health": {"90d": 0.78, "180d": 0.68, "1y": 0.58},
    "maritime_marine": {"90d": 0.80, "180d": 0.72, "1y": 0.62},
    "military_recruitment": {"90d": 0.92, "180d": 0.88, "1y": 0.84},
}

# Industry benchmark QoH scores (fallback)
_INDUSTRY_QOH_BENCHMARKS = {
    "tech_engineering": {"avg_qoh": 68, "top_quartile": 82},
    "healthcare_medical": {"avg_qoh": 65, "top_quartile": 79},
    "finance_banking": {"avg_qoh": 70, "top_quartile": 84},
    "retail_consumer": {"avg_qoh": 55, "top_quartile": 70},
    "logistics_supply_chain": {"avg_qoh": 52, "top_quartile": 68},
    "hospitality_travel": {"avg_qoh": 50, "top_quartile": 66},
    "general_entry_level": {"avg_qoh": 58, "top_quartile": 73},
    "blue_collar_trades": {"avg_qoh": 56, "top_quartile": 71},
    "aerospace_defense": {"avg_qoh": 72, "top_quartile": 86},
    "construction_real_estate": {"avg_qoh": 57, "top_quartile": 72},
    "automotive": {"avg_qoh": 60, "top_quartile": 75},
    "energy_utilities": {"avg_qoh": 68, "top_quartile": 82},
    "pharma_biotech": {"avg_qoh": 70, "top_quartile": 85},
    "insurance": {"avg_qoh": 65, "top_quartile": 79},
    "telecommunications": {"avg_qoh": 62, "top_quartile": 77},
    "food_beverage": {"avg_qoh": 50, "top_quartile": 66},
    "education": {"avg_qoh": 66, "top_quartile": 80},
    "media_entertainment": {"avg_qoh": 60, "top_quartile": 75},
    "legal_services": {"avg_qoh": 68, "top_quartile": 82},
    "mental_health": {"avg_qoh": 60, "top_quartile": 75},
    "maritime_marine": {"avg_qoh": 62, "top_quartile": 77},
    "military_recruitment": {"avg_qoh": 74, "top_quartile": 88},
}

# Source-channel quality modifiers (based on research showing referral/career site
# hires tend to outperform job board hires)
_SOURCE_QUALITY_MODIFIERS = {
    "employee referral": 1.20,
    "referral": 1.20,
    "career site": 1.12,
    "careers page": 1.12,
    "internal": 1.15,
    "linkedin": 1.08,
    "indeed": 0.95,
    "ziprecruiter": 0.93,
    "glassdoor": 0.97,
    "google": 1.00,
    "facebook": 0.92,
    "meta": 0.92,
    "instagram": 0.90,
    "job board": 0.94,
    "programmatic": 0.96,
    "staffing agency": 0.88,
    "agency": 0.88,
    "recruitment agency": 0.88,
    "university": 1.05,
    "campus": 1.05,
    "social media": 0.91,
}


# =============================================================================
# COLUMN MATCHING -- flexible, case-insensitive, partial match
# =============================================================================

_HIRE_COLUMN_PATTERNS: Dict[str, List[str]] = {
    "source": [
        "source",
        "channel",
        "platform",
        "medium",
        "hiring source",
        "recruitment channel",
        "referral source",
        "origin",
    ],
    "hire_date": [
        "hire date",
        "start date",
        "onboard date",
        "date hired",
        "joining date",
        "date of hire",
        "doh",
    ],
    "employee_name": ["employee", "name", "hire name", "candidate", "full name"],
    "employee_id": ["employee id", "emp id", "id", "employee number", "emp no"],
    "role": ["role", "position", "title", "job title", "designation"],
    "department": ["department", "dept", "team", "business unit", "division"],
    "performance_rating": [
        "performance",
        "rating",
        "review score",
        "perf score",
        "performance score",
        "annual review",
        "performance rating",
    ],
    "retention_status": [
        "retention",
        "status",
        "active",
        "retained",
        "still employed",
        "current status",
        "employment status",
        "terminated",
    ],
    "termination_date": [
        "termination date",
        "end date",
        "separation date",
        "exit date",
        "last day",
        "leave date",
        "departure date",
    ],
    "time_to_productivity": [
        "time to productivity",
        "ramp time",
        "ramp up",
        "ttp",
        "productivity days",
        "days to productive",
        "onboarding time",
    ],
    "hiring_manager_satisfaction": [
        "hiring manager",
        "manager satisfaction",
        "hm satisfaction",
        "hm rating",
        "manager rating",
        "hm score",
        "manager score",
    ],
    "cost_to_hire": [
        "cost to hire",
        "cost per hire",
        "cph",
        "hiring cost",
        "recruitment cost",
        "cost",
        "spend",
    ],
    "applications": [
        "applications",
        "applicants",
        "applies",
        "total applicants",
        "candidates applied",
    ],
    "interviews": [
        "interviews",
        "screened",
        "phone screens",
        "interviewed",
        "interview count",
    ],
    "offers": ["offers", "offers made", "offer extended", "offer count"],
    "cultural_fit": [
        "cultural fit",
        "culture score",
        "culture",
        "fit score",
        "value alignment",
    ],
    "diversity_flag": [
        "diversity",
        "diverse",
        "underrepresented",
        "urg",
        "eeo",
        "diversity flag",
        "dei",
    ],
    "gender": ["gender", "sex"],
    "ethnicity": ["ethnicity", "race", "ethnic group"],
}


def _match_column(header: str, patterns: List[str]) -> bool:
    """Check if a column header matches any of the given patterns (case-insensitive, partial)."""
    h = header.lower().strip()
    for pat in patterns:
        if pat in h or h in pat:
            return True
    return False


def _map_columns(headers: List[str]) -> Dict[str, Optional[int]]:
    """Map canonical field names to column indices from the header row."""
    mapping: Dict[str, Optional[int]] = {k: None for k in _HIRE_COLUMN_PATTERNS}
    for idx, header in enumerate(headers):
        if not header or not header.strip():
            continue
        for field, patterns in _HIRE_COLUMN_PATTERNS.items():
            if mapping[field] is None and _match_column(header, patterns):
                mapping[field] = idx
                break
    return mapping


def _safe_float(val: Any) -> Optional[float]:
    """Safely convert a value to float, stripping currency/percent symbols."""
    if val is None:
        return None
    if isinstance(val, (int, float)):
        return float(val)
    s = str(val).strip()
    if not s or s.lower() in ("n/a", "na", "-", "--", "null", "none", ""):
        return None
    s = re.sub(r"[$,\s%]", "", s)
    try:
        return float(s)
    except (ValueError, TypeError):
        return None


def _safe_str(val: Any) -> str:
    """Safely convert a value to trimmed string."""
    if val is None:
        return ""
    return str(val).strip()


def _parse_date(val: Any) -> Optional[datetime.date]:
    """Parse a date from various formats."""
    if val is None:
        return None
    if isinstance(val, datetime.datetime):
        return val.date()
    if isinstance(val, datetime.date):
        return val
    s = str(val).strip()
    if not s or s.lower() in ("n/a", "na", "-", "--", "null", "none"):
        return None
    for fmt in (
        "%Y-%m-%d",
        "%m/%d/%Y",
        "%d/%m/%Y",
        "%m-%d-%Y",
        "%d-%m-%Y",
        "%Y/%m/%d",
        "%B %d, %Y",
        "%b %d, %Y",
        "%d %B %Y",
        "%d %b %Y",
        "%m/%d/%y",
        "%d/%m/%y",
    ):
        try:
            return datetime.datetime.strptime(s, fmt).date()
        except ValueError:
            continue
    return None


def _parse_retention_status(val: Any) -> Optional[bool]:
    """Parse retention status: True = still retained, False = terminated/left."""
    if val is None:
        return None
    s = str(val).strip().lower()
    if s in (
        "active",
        "retained",
        "yes",
        "1",
        "true",
        "current",
        "employed",
        "still employed",
    ):
        return True
    if s in (
        "terminated",
        "left",
        "no",
        "0",
        "false",
        "inactive",
        "departed",
        "resigned",
        "fired",
        "separated",
        "exited",
    ):
        return False
    return None


# =============================================================================
# 1. PARSE HIRING DATA
# =============================================================================


def parse_hiring_data(file_bytes: bytes, filename: str) -> List[Dict[str, Any]]:
    """Parse Excel/CSV with hiring quality data.

    Expected columns (flexible matching):
        Source/Channel, Hire Date, Employee Name, Role, Department,
        Performance Rating (1-5 or 1-10 or 1-100), Retention Status,
        Termination Date, Time to Productivity (days), Hiring Manager
        Satisfaction (1-5 or 1-10), Cost to Hire, Applications,
        Interviews, Offers, Cultural Fit, Diversity Flag

    Returns list of dicts, one per hire record with parsed values.
    """
    try:
        ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
        if ext in ("xlsx", "xls"):
            return _parse_excel_hires(file_bytes)
        elif ext == "csv":
            return _parse_csv_hires(file_bytes)
        else:
            logger.warning("Unsupported file type for hiring data: %s", filename)
            return []
    except Exception as exc:
        logger.exception("Failed to parse hiring data from %s: %s", filename, exc)
        return []


def _parse_excel_hires(file_bytes: bytes) -> List[Dict[str, Any]]:
    """Parse Excel file into hire records."""
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    ws = wb.active
    rows = list(ws.iter_rows(values_only=True))
    wb.close()
    if len(rows) < 2:
        return []
    headers = [str(c).strip() if c else "" for c in rows[0]]
    return _rows_to_hire_records(headers, rows[1:])


def _parse_csv_hires(file_bytes: bytes) -> List[Dict[str, Any]]:
    """Parse CSV file into hire records."""
    import csv as csv_mod

    text = None
    for enc in ("utf-8", "utf-8-sig", "latin-1", "cp1252"):
        try:
            text = file_bytes.decode(enc)
            break
        except (UnicodeDecodeError, ValueError):
            continue
    if not text:
        return []
    reader = csv_mod.reader(io.StringIO(text))
    rows = list(reader)
    if len(rows) < 2:
        return []
    headers = [str(c).strip() if c else "" for c in rows[0]]
    return _rows_to_hire_records(headers, rows[1:])


def _rows_to_hire_records(headers: List[str], data_rows: List) -> List[Dict[str, Any]]:
    """Convert parsed rows to list of hire record dicts."""
    col_map = _map_columns(headers)
    records = []

    for row in data_rows:
        cells = list(row)
        # Skip empty rows
        if not any(c for c in cells if c is not None and str(c).strip()):
            continue

        # Must have at least a source
        src_idx = col_map.get("source")
        if src_idx is None or src_idx >= len(cells) or not cells[src_idx]:
            continue
        source = _safe_str(cells[src_idx])
        if not source or source.lower() in ("total", "grand total", "sum"):
            continue

        def _get_str(field: str) -> str:
            idx = col_map.get(field)
            if idx is not None and idx < len(cells):
                return _safe_str(cells[idx])
            return ""

        def _get_float(field: str) -> Optional[float]:
            idx = col_map.get(field)
            if idx is not None and idx < len(cells):
                return _safe_float(cells[idx])
            return None

        def _get_date(field: str) -> Optional[datetime.date]:
            idx = col_map.get(field)
            if idx is not None and idx < len(cells):
                return _parse_date(cells[idx])
            return None

        def _get_bool(field: str) -> Optional[bool]:
            idx = col_map.get(field)
            if idx is not None and idx < len(cells):
                return _parse_retention_status(cells[idx])
            return None

        # Normalize performance rating to 0-100 scale
        raw_perf = _get_float("performance_rating")
        perf_rating = _normalize_rating(raw_perf)

        # Normalize HM satisfaction to 0-100 scale
        raw_hm = _get_float("hiring_manager_satisfaction")
        hm_satisfaction = _normalize_rating(raw_hm)

        # Normalize cultural fit to 0-100 scale
        raw_culture = _get_float("cultural_fit")
        cultural_fit = _normalize_rating(raw_culture)

        hire_date = _get_date("hire_date")
        term_date = _get_date("termination_date")
        retained = _get_bool("retention_status")

        # Infer retention if we have dates but no explicit status
        if retained is None and hire_date and term_date:
            retained = False
        elif retained is None and hire_date and term_date is None:
            retained = True

        # Calculate tenure in days if possible
        tenure_days = None
        if hire_date:
            end = term_date if term_date else datetime.date.today()
            tenure_days = max(0, (end - hire_date).days)

        rec: Dict[str, Any] = {
            "source": source,
            "hire_date": hire_date.isoformat() if hire_date else None,
            "employee_name": _get_str("employee_name"),
            "employee_id": _get_str("employee_id"),
            "role": _get_str("role"),
            "department": _get_str("department"),
            "performance_rating": perf_rating,
            "retention_status": retained,
            "termination_date": term_date.isoformat() if term_date else None,
            "tenure_days": tenure_days,
            "time_to_productivity": _get_float("time_to_productivity"),
            "hiring_manager_satisfaction": hm_satisfaction,
            "cost_to_hire": _get_float("cost_to_hire") or 0.0,
            "applications": _get_float("applications") or 0.0,
            "interviews": _get_float("interviews") or 0.0,
            "offers": _get_float("offers") or 0.0,
            "cultural_fit": cultural_fit,
            "diversity_flag": _parse_diversity_flag(_get_str("diversity_flag")),
            "gender": _get_str("gender"),
            "ethnicity": _get_str("ethnicity"),
        }
        records.append(rec)

    return records


def _normalize_rating(val: Optional[float]) -> Optional[float]:
    """Normalize a rating to 0-100 scale. Handles 1-5, 1-10, 0-100 scales."""
    if val is None:
        return None
    if val <= 5.0:
        return round(val / 5.0 * 100, 1)
    elif val <= 10.0:
        return round(val / 10.0 * 100, 1)
    else:
        return round(min(100.0, max(0.0, val)), 1)


def _parse_diversity_flag(val: str) -> Optional[bool]:
    """Parse a diversity flag value."""
    if not val:
        return None
    v = val.lower().strip()
    if v in ("yes", "1", "true", "y", "diverse", "underrepresented", "urg"):
        return True
    if v in ("no", "0", "false", "n"):
        return False
    return None


# =============================================================================
# 2. CALCULATE QUALITY OF HIRE (QoH) SCORE
# =============================================================================


def calculate_qoh_score(hire_data: Dict[str, Any] | str | None) -> float:
    """Calculate Quality of Hire composite score (0-100) for a single hire.

    Weighted composite of:
      - Performance rating (30%)
      - 90-day retention (20%)
      - 180-day retention (10%)
      - Time to productivity (15%)
      - Hiring manager satisfaction (15%)
      - Cultural fit (10%)
    """
    # Guard: hire_data may be a string or None instead of dict
    if not isinstance(hire_data, dict):
        logger.warning(
            "calculate_qoh_score received non-dict: %s", type(hire_data).__name__
        )
        return 0.0

    scores = {}
    total_weight = 0.0

    # Performance rating (already 0-100)
    perf = hire_data.get("performance_rating")
    if perf is not None:
        scores["performance_rating"] = perf
        total_weight += _QOH_WEIGHTS["performance_rating"]

    # 90-day retention
    tenure = hire_data.get("tenure_days")
    retained = hire_data.get("retention_status")
    if tenure is not None:
        if tenure >= 90 or retained is True:
            scores["retention_90d"] = 100.0
        elif retained is False and tenure < 90:
            scores["retention_90d"] = max(0.0, tenure / 90.0 * 100.0)
        else:
            scores["retention_90d"] = 100.0  # Still employed, assume retained
        total_weight += _QOH_WEIGHTS["retention_90d"]

        # 180-day retention
        if tenure >= 180 or retained is True:
            scores["retention_180d"] = 100.0
        elif retained is False and tenure < 180:
            scores["retention_180d"] = max(0.0, tenure / 180.0 * 100.0)
        else:
            scores["retention_180d"] = 100.0
        total_weight += _QOH_WEIGHTS["retention_180d"]

    # Time to productivity (lower is better; benchmark 60 days = 100, 180 days = 0)
    ttp = hire_data.get("time_to_productivity")
    if ttp is not None:
        ttp_score = max(0.0, min(100.0, (180.0 - ttp) / 120.0 * 100.0))
        scores["time_to_productivity"] = ttp_score
        total_weight += _QOH_WEIGHTS["time_to_productivity"]

    # Hiring manager satisfaction (already 0-100)
    hm = hire_data.get("hiring_manager_satisfaction")
    if hm is not None:
        scores["hiring_manager_satisfaction"] = hm
        total_weight += _QOH_WEIGHTS["hiring_manager_satisfaction"]

    # Cultural fit (already 0-100)
    cf = hire_data.get("cultural_fit")
    if cf is not None:
        scores["cultural_fit"] = cf
        total_weight += _QOH_WEIGHTS["cultural_fit"]

    if total_weight == 0:
        return 50.0  # No data -- neutral score

    # Weighted average
    weighted_sum = 0.0
    for metric, score in scores.items():
        weighted_sum += score * (_QOH_WEIGHTS[metric] / total_weight)

    return round(min(100.0, max(0.0, weighted_sum)), 1)


def _score_to_grade(score: float) -> str:
    """Convert a 0-100 score to letter grade."""
    if score >= 85:
        return "A"
    elif score >= 70:
        return "B"
    elif score >= 55:
        return "C"
    elif score >= 40:
        return "D"
    else:
        return "F"


# =============================================================================
# 3. ANALYZE SOURCE EFFECTIVENESS
# =============================================================================


def analyze_source_effectiveness(hires: List[Dict[str, Any]]) -> Dict[str, Any]:
    """Full funnel analysis per source: applications -> interviews -> offers -> hires -> retained.

    Returns dict keyed by source with funnel metrics, conversion rates, and QoH scores.
    """
    if not hires:
        return {}

    sources: Dict[str, List[Dict]] = {}
    for h in hires:
        src = h.get("source", "Unknown")
        sources.setdefault(src, []).append(h)

    results: Dict[str, Any] = {}
    all_qoh_scores = []

    for source, source_hires in sources.items():
        n = len(source_hires)

        # Aggregate funnel metrics
        total_apps = sum(h.get("applications") or 0 for h in source_hires)
        total_interviews = sum(h.get("interviews") or 0 for h in source_hires)
        total_offers = sum(h.get("offers") or 0 for h in source_hires)
        total_cost = sum(h.get("cost_to_hire") or 0 for h in source_hires)

        # QoH scores
        qoh_scores = [calculate_qoh_score(h) for h in source_hires]
        avg_qoh = sum(qoh_scores) / len(qoh_scores) if qoh_scores else 50.0
        all_qoh_scores.extend(qoh_scores)

        # Performance ratings
        perf_ratings = [
            h["performance_rating"]
            for h in source_hires
            if h.get("performance_rating") is not None
        ]
        avg_perf = sum(perf_ratings) / len(perf_ratings) if perf_ratings else None

        # Retention
        retained_count = sum(
            1 for h in source_hires if h.get("retention_status") is True
        )
        terminated_count = sum(
            1 for h in source_hires if h.get("retention_status") is False
        )
        known_retention = retained_count + terminated_count
        retention_rate = (
            (retained_count / known_retention * 100) if known_retention > 0 else None
        )

        # Time to productivity
        ttp_values = [
            h["time_to_productivity"]
            for h in source_hires
            if h.get("time_to_productivity") is not None
        ]
        avg_ttp = sum(ttp_values) / len(ttp_values) if ttp_values else None

        # HM satisfaction
        hm_values = [
            h["hiring_manager_satisfaction"]
            for h in source_hires
            if h.get("hiring_manager_satisfaction") is not None
        ]
        avg_hm = sum(hm_values) / len(hm_values) if hm_values else None

        # Conversion rates
        app_to_interview = (
            (total_interviews / total_apps * 100) if total_apps > 0 else 0.0
        )
        interview_to_offer = (
            (total_offers / total_interviews * 100) if total_interviews > 0 else 0.0
        )
        offer_to_hire = (n / total_offers * 100) if total_offers > 0 else 0.0
        app_to_hire = (n / total_apps * 100) if total_apps > 0 else 0.0

        # Cost metrics
        cph = total_cost / n if n > 0 else 0.0
        # CPQH = cost per quality hire (only count hires with QoH >= 60 as "quality")
        quality_hires = sum(1 for q in qoh_scores if q >= 60)
        cpqh = (
            total_cost / quality_hires
            if quality_hires > 0
            else (total_cost if total_cost > 0 else 0.0)
        )
        cpa = total_cost / total_apps if total_apps > 0 else 0.0

        # Diversity
        diverse_count = sum(1 for h in source_hires if h.get("diversity_flag") is True)
        known_diversity = sum(
            1 for h in source_hires if h.get("diversity_flag") is not None
        )
        diversity_rate = (
            (diverse_count / known_diversity * 100) if known_diversity > 0 else None
        )

        results[source] = {
            "source": source,
            "total_hires": n,
            "total_applications": total_apps,
            "total_interviews": total_interviews,
            "total_offers": total_offers,
            "total_cost": round(total_cost, 2),
            "avg_qoh_score": round(avg_qoh, 1),
            "qoh_grade": _score_to_grade(avg_qoh),
            "avg_performance_rating": (
                round(avg_perf, 1) if avg_perf is not None else None
            ),
            "retention_rate": (
                round(retention_rate, 1) if retention_rate is not None else None
            ),
            "retained_count": retained_count,
            "terminated_count": terminated_count,
            "avg_time_to_productivity": (
                round(avg_ttp, 1) if avg_ttp is not None else None
            ),
            "avg_hm_satisfaction": round(avg_hm, 1) if avg_hm is not None else None,
            "quality_hire_count": quality_hires,
            "quality_hire_pct": round(quality_hires / n * 100, 1) if n > 0 else 0.0,
            "cost_per_hire": round(cph, 2),
            "cost_per_quality_hire": round(cpqh, 2),
            "cost_per_application": round(cpa, 2),
            "conversion_rates": {
                "app_to_interview": round(app_to_interview, 1),
                "interview_to_offer": round(interview_to_offer, 1),
                "offer_to_hire": round(offer_to_hire, 1),
                "app_to_hire": round(app_to_hire, 2),
            },
            "diversity_rate": (
                round(diversity_rate, 1) if diversity_rate is not None else None
            ),
            "diverse_hires": diverse_count,
        }

    # Calculate overall stats
    overall_qoh = sum(all_qoh_scores) / len(all_qoh_scores) if all_qoh_scores else 50.0

    return {
        "sources": results,
        "overall_qoh": round(overall_qoh, 1),
        "overall_grade": _score_to_grade(overall_qoh),
        "source_count": len(results),
        "total_hires": len(hires),
    }


# =============================================================================
# 4. CALCULATE RETENTION RATES
# =============================================================================


def calculate_retention_rates(hires: List[Dict[str, Any]]) -> Dict[str, Any]:
    """Calculate 90-day, 180-day, and 1-year retention rates by source.

    Returns dict with retention data per source and overall.
    """
    if not hires:
        return {"sources": {}, "overall": {}}

    sources: Dict[str, List[Dict]] = {}
    for h in hires:
        src = h.get("source", "Unknown")
        sources.setdefault(src, []).append(h)

    source_retention: Dict[str, Dict] = {}

    for source, source_hires in sources.items():
        rates = _compute_retention_for_group(source_hires)
        source_retention[source] = rates

    overall_rates = _compute_retention_for_group(hires)

    return {
        "sources": source_retention,
        "overall": overall_rates,
    }


def _compute_retention_for_group(group: List[Dict]) -> Dict[str, Any]:
    """Compute retention rates for a group of hires."""
    # For each milestone, count hires who had enough time to reach it
    milestones = {"90d": 90, "180d": 180, "1y": 365}
    rates: Dict[str, Any] = {}

    for label, days in milestones.items():
        eligible = 0
        retained = 0

        for h in group:
            hire_date = (
                _parse_date(h.get("hire_date"))
                if isinstance(h.get("hire_date"), str)
                else h.get("hire_date")
            )
            if hire_date is None:
                continue

            # Check if enough time has passed since hire date
            if isinstance(hire_date, str):
                try:
                    hire_date = datetime.date.fromisoformat(hire_date)
                except (ValueError, TypeError):
                    continue

            days_since_hire = (
                (datetime.date.today() - hire_date).days
                if isinstance(hire_date, datetime.date)
                else 0
            )
            if days_since_hire < days:
                continue  # Not enough time has passed

            eligible += 1

            tenure = h.get("tenure_days")
            status = h.get("retention_status")

            if status is True:
                retained += 1
            elif status is False and tenure is not None and tenure >= days:
                retained += 1
            elif status is False:
                pass  # Left before milestone
            elif tenure is not None and tenure >= days:
                retained += 1

        rate = (retained / eligible * 100) if eligible > 0 else None
        rates[label] = {
            "rate": round(rate, 1) if rate is not None else None,
            "eligible": eligible,
            "retained": retained,
        }

    return rates


# =============================================================================
# 5. CALCULATE COST PER QUALITY HIRE (CPQH)
# =============================================================================


def calculate_cost_per_quality_hire(hires: List[Dict[str, Any]]) -> Dict[str, Any]:
    """Calculate CPQH metric -- cost per quality hire vs traditional cost per hire.

    A "quality hire" is defined as one with QoH score >= 60.
    """
    if not hires:
        return {"sources": {}, "overall": {}}

    sources: Dict[str, List[Dict]] = {}
    for h in hires:
        src = h.get("source", "Unknown")
        sources.setdefault(src, []).append(h)

    source_costs: Dict[str, Dict] = {}
    total_cost_all = 0.0
    total_hires_all = 0
    total_quality_all = 0

    for source, source_hires in sources.items():
        n = len(source_hires)
        total_cost = sum(h.get("cost_to_hire") or 0 for h in source_hires)
        qoh_scores = [calculate_qoh_score(h) for h in source_hires]
        quality_count = sum(1 for q in qoh_scores if q >= 60)

        cph = total_cost / n if n > 0 else 0.0
        cpqh = total_cost / quality_count if quality_count > 0 else None
        cpa = (
            total_cost / sum(h.get("applications") or 0 for h in source_hires)
            if sum(h.get("applications") or 0 for h in source_hires) > 0
            else 0.0
        )

        # Value ratio: how much more expensive is a quality hire vs any hire?
        value_ratio = (cpqh / cph) if (cpqh and cph > 0) else None

        source_costs[source] = {
            "source": source,
            "total_cost": round(total_cost, 2),
            "total_hires": n,
            "quality_hires": quality_count,
            "quality_pct": round(quality_count / n * 100, 1) if n > 0 else 0.0,
            "cost_per_hire": round(cph, 2),
            "cost_per_quality_hire": round(cpqh, 2) if cpqh is not None else None,
            "cost_per_application": round(cpa, 2),
            "value_ratio": round(value_ratio, 2) if value_ratio is not None else None,
            "is_cost_effective": (value_ratio is not None and value_ratio < 1.5),
        }

        total_cost_all += total_cost
        total_hires_all += n
        total_quality_all += quality_count

    overall_cph = total_cost_all / total_hires_all if total_hires_all > 0 else 0.0
    overall_cpqh = total_cost_all / total_quality_all if total_quality_all > 0 else None

    return {
        "sources": source_costs,
        "overall": {
            "total_cost": round(total_cost_all, 2),
            "total_hires": total_hires_all,
            "total_quality_hires": total_quality_all,
            "quality_pct": (
                round(total_quality_all / total_hires_all * 100, 1)
                if total_hires_all > 0
                else 0.0
            ),
            "cost_per_hire": round(overall_cph, 2),
            "cost_per_quality_hire": (
                round(overall_cpqh, 2) if overall_cpqh is not None else None
            ),
        },
    }


# =============================================================================
# 6. BENCHMARK AGAINST INDUSTRY
# =============================================================================


def benchmark_against_industry(
    results: Dict[str, Any], industry: str = "general_entry_level"
) -> Dict[str, Any]:
    """Compare hiring quality metrics against industry benchmarks.

    Returns dict with benchmark comparisons and variance analysis.
    """
    bench_retention = _INDUSTRY_RETENTION_BENCHMARKS.get(
        industry, _INDUSTRY_RETENTION_BENCHMARKS["general_entry_level"]
    )
    bench_qoh = _INDUSTRY_QOH_BENCHMARKS.get(
        industry, _INDUSTRY_QOH_BENCHMARKS["general_entry_level"]
    )

    overall_qoh = results.get("overall_qoh", 50.0)
    avg_benchmark_qoh = bench_qoh["avg_qoh"]
    top_quartile_qoh = bench_qoh["top_quartile"]

    qoh_vs_avg = overall_qoh - avg_benchmark_qoh
    qoh_vs_top = overall_qoh - top_quartile_qoh

    # Retention benchmarks
    retention_data = results.get("retention", {}).get("overall", {})
    retention_comparisons = {}
    for period, days_label in [("90d", "90d"), ("180d", "180d"), ("1y", "1y")]:
        actual_data = retention_data.get(period, {})
        actual_rate = actual_data.get("rate")
        bench_rate = bench_retention.get(period, 0.5) * 100
        if actual_rate is not None:
            variance = actual_rate - bench_rate
            retention_comparisons[period] = {
                "actual": actual_rate,
                "benchmark": round(bench_rate, 1),
                "variance": round(variance, 1),
                "is_above_benchmark": variance >= 0,
            }
        else:
            retention_comparisons[period] = {
                "actual": None,
                "benchmark": round(bench_rate, 1),
                "variance": None,
                "is_above_benchmark": None,
            }

    # Determine overall performance tier
    if qoh_vs_avg >= 10:
        tier = "top_performer"
        tier_label = "Top Performer"
        tier_description = (
            "Your quality of hire significantly exceeds industry average."
        )
    elif qoh_vs_avg >= 0:
        tier = "above_average"
        tier_label = "Above Average"
        tier_description = "Your quality of hire is above the industry average."
    elif qoh_vs_avg >= -10:
        tier = "near_average"
        tier_label = "Near Average"
        tier_description = "Your quality of hire is close to the industry average."
    else:
        tier = "below_average"
        tier_label = "Below Average"
        tier_description = "Your quality of hire falls below the industry average. Focus on top-performing sources."

    return {
        "industry": industry,
        "industry_label": INDUSTRY_LABEL_MAP.get(
            industry, industry.replace("_", " ").title()
        ),
        "qoh_comparison": {
            "actual": overall_qoh,
            "industry_avg": avg_benchmark_qoh,
            "top_quartile": top_quartile_qoh,
            "vs_avg": round(qoh_vs_avg, 1),
            "vs_top_quartile": round(qoh_vs_top, 1),
            "is_above_avg": qoh_vs_avg >= 0,
            "is_top_quartile": qoh_vs_top >= 0,
        },
        "retention_comparison": retention_comparisons,
        "tier": tier,
        "tier_label": tier_label,
        "tier_description": tier_description,
    }


# =============================================================================
# 7. GENERATE RECOMMENDATIONS
# =============================================================================


def generate_recommendations(analysis: Dict[str, Any]) -> List[Dict[str, Any]]:
    """Generate actionable recommendations based on the analysis.

    Returns list of recommendation dicts ranked by impact score.
    """
    recommendations = []

    source_data = analysis.get("source_effectiveness", {}).get("sources", {})
    retention_data = analysis.get("retention", {})
    cost_data = analysis.get("cost_analysis", {}).get("sources", {})
    benchmark_data = analysis.get("benchmark", {})

    if not source_data:
        return [
            {
                "title": "Insufficient Data",
                "description": "Upload hiring data with source/channel information to receive recommendations.",
                "impact_score": 0,
                "category": "data",
                "priority": "medium",
            }
        ]

    # Sort sources by QoH
    sorted_sources = sorted(
        source_data.values(), key=lambda x: x.get("avg_qoh_score") or 0, reverse=True
    )

    # 1. Scale top-performing sources
    for src in sorted_sources:
        if src["avg_qoh_score"] >= 70 and src["total_hires"] >= 2:
            impact = min(95, int(src["avg_qoh_score"] * 0.9 + src["total_hires"] * 2))
            recommendations.append(
                {
                    "title": f"Scale {src['source']} -- Top Quality Source",
                    "description": (
                        f"{src['source']} produces hires with avg QoH of {src['avg_qoh_score']}/100 (Grade {src['qoh_grade']}). "
                        f"Quality hire rate: {src['quality_hire_pct']:.0f}%. "
                        f"Increase investment in this channel by 20-30% to maximize quality hires."
                    ),
                    "impact_score": impact,
                    "category": "scale",
                    "priority": "high",
                    "source": src["source"],
                }
            )

    # 2. Reduce low-quality sources
    for src in reversed(sorted_sources):
        if src["avg_qoh_score"] < 45 and src["total_hires"] >= 2:
            impact = min(
                90,
                int(
                    (100 - src["avg_qoh_score"]) * 0.8 + src.get("total_cost")
                    or 0 / 1000
                ),
            )
            recommendations.append(
                {
                    "title": f"Reduce {src['source']} -- Low Quality Source",
                    "description": (
                        f"{src['source']} produces hires with avg QoH of only {src['avg_qoh_score']}/100 (Grade {src['qoh_grade']}). "
                        f"Quality hire rate: {src['quality_hire_pct']:.0f}%. "
                        f"Reduce budget by 25-40% and redirect to top-performing channels."
                    ),
                    "impact_score": impact,
                    "category": "reduce",
                    "priority": "high",
                    "source": src["source"],
                }
            )

    # 3. CPQH optimization
    for source_name, cost_info in cost_data.items():
        if cost_info.get("value_ratio") and cost_info["value_ratio"] > 2.0:
            recommendations.append(
                {
                    "title": f"Improve {source_name} Quality Yield",
                    "description": (
                        f"{source_name} has a CPQH/CPH ratio of {cost_info['value_ratio']:.1f}x -- "
                        f"meaning you pay {cost_info['value_ratio']:.1f}x more per quality hire than per any hire. "
                        f"Review screening criteria and candidate matching for this source."
                    ),
                    "impact_score": min(85, int(cost_info["value_ratio"] * 20)),
                    "category": "optimize",
                    "priority": "medium",
                    "source": source_name,
                }
            )

    # 4. Retention-focused recommendations
    for source_name, src_data in source_data.items():
        ret_rate = src_data.get("retention_rate")
        if ret_rate is not None and ret_rate < 60 and src_data["total_hires"] >= 3:
            recommendations.append(
                {
                    "title": f"Address {source_name} Retention Problem",
                    "description": (
                        f"{source_name} has only {ret_rate:.0f}% retention rate. "
                        f"Review job ad accuracy, interview process, and onboarding for hires from this source. "
                        f"Consider pre-hire assessments to improve candidate-role fit."
                    ),
                    "impact_score": min(90, int((100 - ret_rate) * 0.9)),
                    "category": "retention",
                    "priority": "high",
                    "source": source_name,
                }
            )

    # 5. Funnel efficiency
    for source_name, src_data in source_data.items():
        conv = src_data.get("conversion_rates", {})
        app_to_hire = conv.get("app_to_hire") or 0
        if (
            app_to_hire > 0
            and app_to_hire < 1.0
            and src_data.get("total_applications")
            or 0 > 20
        ):
            recommendations.append(
                {
                    "title": f"Improve {source_name} Funnel Efficiency",
                    "description": (
                        f"{source_name} converts only {app_to_hire:.2f}% of applications to hires. "
                        f"Review job descriptions, screening criteria, and interview process. "
                        f"Better targeting can reduce wasted spend on unqualified applicants."
                    ),
                    "impact_score": min(75, int((5.0 - app_to_hire) * 15)),
                    "category": "funnel",
                    "priority": "medium",
                    "source": source_name,
                }
            )

    # 6. Benchmark-based recommendations
    tier = benchmark_data.get("tier", "near_average")
    if tier == "below_average":
        recommendations.append(
            {
                "title": "Quality of Hire Below Industry Average",
                "description": (
                    f"Your overall QoH ({benchmark_data.get('qoh_comparison', {}).get('actual') or 0:.0f}) "
                    f"is below the industry average ({benchmark_data.get('qoh_comparison', {}).get('industry_avg') or 0}). "
                    f"Focus spending on your top 2-3 quality sources and improve screening processes."
                ),
                "impact_score": 88,
                "category": "benchmark",
                "priority": "critical",
            }
        )

    # 7. Diversity recommendations
    diverse_sources = [
        (s, d)
        for s, d in source_data.items()
        if d.get("diversity_rate") is not None and d["diversity_rate"] > 30
    ]
    if diverse_sources:
        best_diverse = max(diverse_sources, key=lambda x: x[1]["diversity_rate"])
        recommendations.append(
            {
                "title": f"Leverage {best_diverse[0]} for Diversity Goals",
                "description": (
                    f"{best_diverse[0]} has a {best_diverse[1]['diversity_rate']:.0f}% diversity rate -- "
                    f"the highest among your sources. Scale this channel to improve overall diversity metrics."
                ),
                "impact_score": 70,
                "category": "diversity",
                "priority": "medium",
                "source": best_diverse[0],
            }
        )

    # Sort by impact score
    recommendations.sort(key=lambda x: x.get("impact_score") or 0, reverse=True)

    return recommendations[:12]  # Cap at 12


# =============================================================================
# 8. PREDICT SOURCE PERFORMANCE
# =============================================================================


def predict_source_performance(analysis: Dict[str, Any]) -> Dict[str, Any]:
    """ML-lite predictions for next quarter based on source patterns.

    Uses weighted moving trends and source quality modifiers to project
    future performance by source.
    """
    source_data = analysis.get("source_effectiveness", {}).get("sources", {})
    if not source_data:
        return {"predictions": {}, "confidence": "low"}

    predictions: Dict[str, Dict] = {}

    for source_name, src in source_data.items():
        current_qoh = src.get("avg_qoh_score", 50.0)
        retention_rate = src.get("retention_rate")
        quality_pct = src.get("quality_hire_pct", 50.0)
        hires = src.get("total_hires") or 0

        # Source quality modifier
        modifier = 1.0
        src_lower = source_name.lower()
        for pattern, mod in _SOURCE_QUALITY_MODIFIERS.items():
            if pattern in src_lower:
                modifier = mod
                break

        # Confidence based on sample size
        if hires >= 20:
            confidence = "high"
            variance_factor = 0.05
        elif hires >= 10:
            confidence = "medium"
            variance_factor = 0.10
        elif hires >= 5:
            confidence = "low"
            variance_factor = 0.15
        else:
            confidence = "very_low"
            variance_factor = 0.25

        # Project QoH (mean reversion toward 60 + source modifier influence)
        projected_qoh = current_qoh * 0.7 + (60 * modifier) * 0.3
        projected_qoh = round(min(100, max(0, projected_qoh)), 1)

        # Project retention (slight regression to industry mean)
        if retention_rate is not None:
            projected_retention = retention_rate * 0.8 + 70 * 0.2
        else:
            projected_retention = 70.0
        projected_retention = round(min(100, max(0, projected_retention)), 1)

        # Project quality hire percentage
        projected_quality = quality_pct * 0.75 + 55 * 0.25
        projected_quality = round(min(100, max(0, projected_quality)), 1)

        # Trend direction
        if projected_qoh > current_qoh + 2:
            trend = "improving"
        elif projected_qoh < current_qoh - 2:
            trend = "declining"
        else:
            trend = "stable"

        predictions[source_name] = {
            "source": source_name,
            "current_qoh": current_qoh,
            "projected_qoh": projected_qoh,
            "projected_grade": _score_to_grade(projected_qoh),
            "projected_retention": projected_retention,
            "projected_quality_pct": projected_quality,
            "trend": trend,
            "confidence": confidence,
            "range_low": round(max(0, projected_qoh * (1 - variance_factor)), 1),
            "range_high": round(min(100, projected_qoh * (1 + variance_factor)), 1),
        }

    # Overall confidence
    total_hires = sum(s.get("total_hires") or 0 for s in source_data.values())
    overall_confidence = (
        "high" if total_hires >= 50 else ("medium" if total_hires >= 20 else "low")
    )

    return {
        "predictions": predictions,
        "confidence": overall_confidence,
        "methodology": "Weighted trend projection with source quality modifiers and mean reversion",
    }


# =============================================================================
# 9. GENERATE SAMPLE DATA
# =============================================================================


def generate_sample_data() -> List[Dict[str, Any]]:
    """Generate realistic sample hiring data for demo purposes.

    Returns list of hire records that can be analyzed.
    """
    sources = [
        ("LinkedIn", 0.35),
        ("Indeed", 0.25),
        ("Employee Referral", 0.15),
        ("Career Site", 0.10),
        ("Glassdoor", 0.05),
        ("University Recruiting", 0.05),
        ("Staffing Agency", 0.05),
    ]

    roles = [
        "Software Engineer",
        "Product Manager",
        "Data Analyst",
        "Marketing Manager",
        "Sales Representative",
        "UX Designer",
        "DevOps Engineer",
        "HR Coordinator",
        "Financial Analyst",
        "Customer Success Manager",
    ]

    departments = [
        "Engineering",
        "Product",
        "Data",
        "Marketing",
        "Sales",
        "Design",
        "Engineering",
        "People Ops",
        "Finance",
        "Customer Success",
    ]

    first_names = [
        "Alex",
        "Jordan",
        "Taylor",
        "Morgan",
        "Casey",
        "Riley",
        "Avery",
        "Quinn",
        "Sage",
        "Reese",
        "Dakota",
        "Cameron",
        "Emerson",
        "Finley",
        "Hayden",
        "Jamie",
        "Kendall",
        "Logan",
        "Parker",
        "Peyton",
        "Robin",
        "Sam",
        "Skyler",
        "Sydney",
    ]
    last_names = [
        "Smith",
        "Johnson",
        "Williams",
        "Brown",
        "Jones",
        "Garcia",
        "Miller",
        "Davis",
        "Rodriguez",
        "Martinez",
        "Hernandez",
        "Lopez",
        "Wilson",
        "Anderson",
        "Thomas",
        "Taylor",
        "Moore",
        "Jackson",
        "Martin",
        "Lee",
        "Perez",
        "Thompson",
        "White",
        "Harris",
    ]

    records = []
    base_date = datetime.date.today() - datetime.timedelta(days=400)

    for i in range(75):
        # Pick source with weighting
        rand_val = random.random()
        cumulative = 0.0
        source = sources[0][0]
        for src_name, weight in sources:
            cumulative += weight
            if rand_val <= cumulative:
                source = src_name
                break

        # Source-specific quality profiles
        src_lower = source.lower()
        if "referral" in src_lower:
            base_perf = random.gauss(82, 8)
            base_retention_prob = 0.88
            base_ttp = random.gauss(35, 10)
            base_hm = random.gauss(85, 7)
            base_cost = random.gauss(2500, 500)
            apps_per_hire = random.randint(3, 8)
        elif "career site" in src_lower:
            base_perf = random.gauss(76, 10)
            base_retention_prob = 0.82
            base_ttp = random.gauss(40, 12)
            base_hm = random.gauss(78, 10)
            base_cost = random.gauss(3000, 800)
            apps_per_hire = random.randint(15, 35)
        elif "linkedin" in src_lower:
            base_perf = random.gauss(72, 12)
            base_retention_prob = 0.78
            base_ttp = random.gauss(45, 12)
            base_hm = random.gauss(74, 10)
            base_cost = random.gauss(4500, 1200)
            apps_per_hire = random.randint(20, 50)
        elif "indeed" in src_lower:
            base_perf = random.gauss(62, 15)
            base_retention_prob = 0.65
            base_ttp = random.gauss(55, 15)
            base_hm = random.gauss(65, 12)
            base_cost = random.gauss(3200, 900)
            apps_per_hire = random.randint(40, 100)
        elif "university" in src_lower:
            base_perf = random.gauss(70, 10)
            base_retention_prob = 0.80
            base_ttp = random.gauss(60, 15)
            base_hm = random.gauss(72, 10)
            base_cost = random.gauss(5000, 1500)
            apps_per_hire = random.randint(25, 60)
        elif "agency" in src_lower:
            base_perf = random.gauss(58, 14)
            base_retention_prob = 0.55
            base_ttp = random.gauss(50, 15)
            base_hm = random.gauss(60, 14)
            base_cost = random.gauss(8000, 2000)
            apps_per_hire = random.randint(5, 15)
        else:  # Glassdoor, etc.
            base_perf = random.gauss(65, 13)
            base_retention_prob = 0.70
            base_ttp = random.gauss(48, 13)
            base_hm = random.gauss(68, 11)
            base_cost = random.gauss(3500, 1000)
            apps_per_hire = random.randint(30, 70)

        hire_date = base_date + datetime.timedelta(days=random.randint(0, 380))
        is_retained = random.random() < base_retention_prob
        tenure_days = (
            (datetime.date.today() - hire_date).days
            if is_retained
            else random.randint(15, min(300, (datetime.date.today() - hire_date).days))
        )
        term_date = (
            None if is_retained else (hire_date + datetime.timedelta(days=tenure_days))
        )

        role_idx = i % len(roles)
        perf = max(1.0, min(5.0, base_perf / 20.0))
        hm_sat = max(1.0, min(5.0, base_hm / 20.0))
        cultural = max(1.0, min(5.0, random.gauss(base_perf / 20.0 - 0.2, 0.5)))
        ttp = max(10, int(base_ttp))
        cost = max(500, round(base_cost, 2))

        total_apps = apps_per_hire
        interviews = max(1, int(total_apps * random.uniform(0.15, 0.40)))
        offers = max(1, int(interviews * random.uniform(0.30, 0.60)))

        is_diverse = random.random() < 0.35

        records.append(
            {
                "source": source,
                "hire_date": hire_date.isoformat(),
                "employee_name": f"{random.choice(first_names)} {random.choice(last_names)}",
                "employee_id": f"EMP-{1000 + i}",
                "role": roles[role_idx],
                "department": departments[role_idx],
                "performance_rating": round(perf, 1),
                "retention_status": is_retained,
                "termination_date": term_date.isoformat() if term_date else None,
                "tenure_days": tenure_days,
                "time_to_productivity": ttp,
                "hiring_manager_satisfaction": round(hm_sat, 1),
                "cost_to_hire": cost,
                "applications": total_apps,
                "interviews": interviews,
                "offers": offers,
                "cultural_fit": round(cultural, 1),
                "diversity_flag": is_diverse,
                "gender": "",
                "ethnicity": "",
            }
        )

    return records


# =============================================================================
# 10. GENERATE SIGNAL EXCEL
# =============================================================================


def generate_signal_excel(report: Dict[str, Any], client_name: str = "Client") -> bytes:
    """Generate multi-sheet Excel report with HireSignal analysis.

    Sheets:
        1. Executive Summary
        2. Source Effectiveness Matrix
        3. Retention Analysis
        4. Cost Intelligence
        5. Recommendations
        6. Predictions

    Uses Sapphire Blue palette, Calibri font, data starts at column B.
    """
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
        from openpyxl.utils import get_column_letter
    except ImportError:
        logger.error("openpyxl not available for Excel generation")
        return b""

    wb = Workbook()

    # Design tokens
    f_title = Font(name="Calibri", bold=True, size=18, color=_WHITE)
    f_section = Font(name="Calibri", bold=True, size=14, color=_WHITE)
    f_subsection = Font(name="Calibri", bold=True, size=12, color=_NAVY)
    f_header = Font(name="Calibri", bold=True, size=10, color=_WHITE)
    f_body = Font(name="Calibri", size=10, color=_STONE)
    f_body_bold = Font(name="Calibri", bold=True, size=10, color=_STONE)
    f_footnote = Font(name="Calibri", italic=True, size=9, color=_MUTED)
    f_green = Font(name="Calibri", bold=True, size=10, color=_GREEN)
    f_red = Font(name="Calibri", bold=True, size=10, color=_RED)
    f_amber = Font(name="Calibri", bold=True, size=10, color=_AMBER)

    fill_navy = PatternFill("solid", fgColor=_NAVY)
    fill_sapphire = PatternFill("solid", fgColor=_SAPPHIRE)
    fill_light = PatternFill("solid", fgColor=_BLUE_LIGHT)
    fill_pale = PatternFill("solid", fgColor=_BLUE_PALE)
    fill_white = PatternFill("solid", fgColor=_WHITE)
    fill_green_bg = PatternFill("solid", fgColor=_GREEN_BG)
    fill_red_bg = PatternFill("solid", fgColor=_RED_BG)
    fill_amber_bg = PatternFill("solid", fgColor=_AMBER_BG)

    align_center = Alignment(horizontal="center", vertical="center", wrap_text=True)
    align_left = Alignment(horizontal="left", vertical="center", wrap_text=True)
    align_right = Alignment(horizontal="right", vertical="center")

    thin_border = Border(
        left=Side(style="thin", color=_WARM_GRAY),
        right=Side(style="thin", color=_WARM_GRAY),
        top=Side(style="thin", color=_WARM_GRAY),
        bottom=Side(style="thin", color=_WARM_GRAY),
    )

    B = COL_START

    def _set_col_widths(ws, widths):
        for i, w in enumerate(widths, start=1):
            ws.column_dimensions[get_column_letter(i)].width = w

    def _write_section_header(ws, row, title, col_span=7):
        for c in range(B, B + col_span):
            ws.cell(row=row, column=c).fill = fill_navy
        cell = ws.cell(row=row, column=B, value=title)
        cell.font = f_section
        cell.alignment = align_left
        ws.row_dimensions[row].height = 32
        return row + 1

    def _write_table_header(ws, row, headers):
        for i, h in enumerate(headers):
            cell = ws.cell(row=row, column=B + i, value=h)
            cell.font = f_header
            cell.fill = fill_sapphire
            cell.alignment = align_center
            cell.border = thin_border
        ws.row_dimensions[row].height = 28
        return row + 1

    def _write_table_row(ws, row, values, fonts=None, fills=None, aligns=None):
        for i, v in enumerate(values):
            cell = ws.cell(row=row, column=B + i, value=v)
            cell.font = fonts[i] if fonts and i < len(fonts) else f_body
            cell.fill = fills[i] if fills and i < len(fills) else fill_white
            cell.alignment = aligns[i] if aligns and i < len(aligns) else align_center
            cell.border = thin_border
        ws.row_dimensions[row].height = 24
        return row + 1

    # Extract data
    source_eff = report.get("source_effectiveness", {})
    sources = source_eff.get("sources", {})
    retention = report.get("retention", {})
    cost_analysis = report.get("cost_analysis", {})
    benchmark = report.get("benchmark", {})
    recommendations = report.get("recommendations") or []
    predictions = report.get("predictions", {}).get("predictions", {})
    overall_qoh = source_eff.get("overall_qoh", 50.0)
    overall_grade = source_eff.get("overall_grade", "C")

    # ── SHEET 1: Executive Summary ──────────────────────────────────────
    ws1 = wb.active
    ws1.title = "Executive Summary"
    _set_col_widths(ws1, [3, 22, 18, 18, 18, 18, 18, 18, 3])

    row = 2
    for c in range(B, B + 7):
        ws1.cell(row=row, column=c).fill = fill_navy
    ws1.cell(row=row, column=B, value=f"HireSignal -- Quality of Hire Report").font = (
        f_title
    )
    ws1.cell(row=row, column=B).alignment = align_left
    ws1.row_dimensions[row].height = 40
    row += 1
    ws1.cell(
        row=row,
        column=B,
        value=f"{client_name} | Generated {datetime.datetime.now().strftime('%B %d, %Y')}",
    ).font = Font(name="Calibri", size=11, color=_SAPPHIRE)
    row += 2

    row = _write_section_header(ws1, row, "Quality of Hire Scorecard")

    # Scorecard metrics
    row = _write_table_header(
        ws1,
        row,
        ["Metric", "Value", "Grade", "Industry Avg", "vs. Benchmark", "Status", ""],
    )
    qoh_comp = benchmark.get("qoh_comparison", {})
    tier_label = benchmark.get("tier_label", "N/A")

    scorecard_rows = [
        (
            "Overall QoH Score",
            f"{overall_qoh:.1f}/100",
            overall_grade,
            str(qoh_comp.get("industry_avg", "N/A")),
            f"{qoh_comp.get('vs_avg') or 0:+.1f}",
            tier_label,
            "",
        ),
        (
            "Total Hires Analyzed",
            str(source_eff.get("total_hires") or 0),
            "",
            "",
            "",
            "",
            "",
        ),
        (
            "Sources Tracked",
            str(source_eff.get("source_count") or 0),
            "",
            "",
            "",
            "",
            "",
        ),
    ]

    overall_cost = cost_analysis.get("overall", {})
    if overall_cost:
        scorecard_rows.extend(
            [
                (
                    "Cost per Hire",
                    f"${overall_cost.get('cost_per_hire') or 0:,.2f}",
                    "",
                    "",
                    "",
                    "",
                    "",
                ),
                (
                    "Cost per Quality Hire",
                    (
                        f"${overall_cost.get('cost_per_quality_hire') or 0:,.2f}"
                        if overall_cost.get("cost_per_quality_hire")
                        else "N/A"
                    ),
                    "",
                    "",
                    "",
                    "",
                    "",
                ),
                (
                    "Quality Hire %",
                    f"{overall_cost.get('quality_pct') or 0:.1f}%",
                    "",
                    "",
                    "",
                    "",
                    "",
                ),
            ]
        )

    for sr in scorecard_rows:
        grade_val = sr[2]
        if grade_val in ("A", "B"):
            grade_font = f_green
        elif grade_val in ("D", "F"):
            grade_font = f_red
        elif grade_val == "C":
            grade_font = f_amber
        else:
            grade_font = f_body
        row = _write_table_row(
            ws1,
            row,
            list(sr),
            fonts=[f_body_bold, f_body, grade_font, f_body, f_body, f_body, f_body],
        )

    row += 2
    ws1.cell(row=row, column=B, value="Powered by HireSignal | Joveo Nova AI").font = (
        f_footnote
    )

    # ── SHEET 2: Source Effectiveness Matrix ─────────────────────────────
    ws2 = wb.create_sheet("Source Effectiveness")
    _set_col_widths(ws2, [3, 22, 12, 12, 10, 10, 12, 14, 14, 12, 3])

    row = 2
    row = _write_section_header(ws2, row, "Source Effectiveness Matrix", col_span=9)
    row = _write_table_header(
        ws2,
        row,
        [
            "Source",
            "Hires",
            "Applications",
            "Interviews",
            "Offers",
            "QoH Score",
            "QoH Grade",
            "CPQH",
            "Retention %",
        ],
    )

    sorted_sources = sorted(
        sources.values(), key=lambda x: x.get("avg_qoh_score") or 0, reverse=True
    )
    for src in sorted_sources:
        grade = src.get("qoh_grade", "C")
        grade_font = (
            f_green
            if grade in ("A", "B")
            else (f_red if grade in ("D", "F") else f_amber)
        )
        ret_val = (
            f"{src['retention_rate']:.0f}%"
            if src.get("retention_rate") is not None
            else "N/A"
        )
        cpqh_val = (
            f"${src['cost_per_quality_hire']:,.0f}"
            if src.get("cost_per_quality_hire")
            else "N/A"
        )

        row = _write_table_row(
            ws2,
            row,
            [
                src["source"],
                src["total_hires"],
                int(src.get("total_applications") or 0),
                int(src.get("total_interviews") or 0),
                int(src.get("total_offers") or 0),
                f"{src['avg_qoh_score']:.1f}",
                grade,
                cpqh_val,
                ret_val,
            ],
            fonts=[
                f_body_bold,
                f_body,
                f_body,
                f_body,
                f_body,
                f_body_bold,
                grade_font,
                f_body,
                f_body,
            ],
        )

    row += 1
    # Conversion rates sub-table
    row = _write_section_header(ws2, row, "Funnel Conversion Rates", col_span=9)
    row = _write_table_header(
        ws2,
        row,
        [
            "Source",
            "App->Interview",
            "Interview->Offer",
            "Offer->Hire",
            "App->Hire",
            "",
            "",
            "",
            "",
        ],
    )
    for src in sorted_sources:
        conv = src.get("conversion_rates", {})
        row = _write_table_row(
            ws2,
            row,
            [
                src["source"],
                f"{conv.get('app_to_interview') or 0:.1f}%",
                f"{conv.get('interview_to_offer') or 0:.1f}%",
                f"{conv.get('offer_to_hire') or 0:.1f}%",
                f"{conv.get('app_to_hire') or 0:.2f}%",
                "",
                "",
                "",
                "",
            ],
            fonts=[f_body_bold] + [f_body] * 8,
        )

    # ── SHEET 3: Retention Analysis ──────────────────────────────────────
    ws3 = wb.create_sheet("Retention Analysis")
    _set_col_widths(ws3, [3, 22, 14, 14, 14, 14, 14, 3])

    row = 2
    row = _write_section_header(ws3, row, "Retention by Source", col_span=6)
    row = _write_table_header(
        ws3,
        row,
        [
            "Source",
            "90-Day Rate",
            "180-Day Rate",
            "1-Year Rate",
            "Retained",
            "Terminated",
        ],
    )

    source_retention = retention.get("sources", {})
    for source_name in sorted(source_retention.keys()):
        sr = source_retention[source_name]
        r90 = sr.get("90d", {}).get("rate")
        r180 = sr.get("180d", {}).get("rate")
        r1y = sr.get("1y", {}).get("rate")
        src_info = sources.get(source_name, {})

        row = _write_table_row(
            ws3,
            row,
            [
                source_name,
                f"{r90:.0f}%" if r90 is not None else "N/A",
                f"{r180:.0f}%" if r180 is not None else "N/A",
                f"{r1y:.0f}%" if r1y is not None else "N/A",
                src_info.get("retained_count") or 0,
                src_info.get("terminated_count") or 0,
            ],
            fonts=[f_body_bold] + [f_body] * 5,
        )

    # Overall row
    overall_ret = retention.get("overall", {})
    row = _write_table_row(
        ws3,
        row,
        [
            "OVERALL",
            (
                f"{overall_ret.get('90d', {}).get('rate') or 0:.0f}%"
                if overall_ret.get("90d", {}).get("rate")
                else "N/A"
            ),
            (
                f"{overall_ret.get('180d', {}).get('rate') or 0:.0f}%"
                if overall_ret.get("180d", {}).get("rate")
                else "N/A"
            ),
            (
                f"{overall_ret.get('1y', {}).get('rate') or 0:.0f}%"
                if overall_ret.get("1y", {}).get("rate")
                else "N/A"
            ),
            "",
            "",
        ],
        fonts=[f_body_bold] * 2 + [f_body_bold] * 4,
        fills=[fill_light] * 6,
    )

    # Benchmark comparison
    ret_comp = benchmark.get("retention_comparison", {})
    if ret_comp:
        row += 1
        row = _write_section_header(
            ws3, row, "Retention vs Industry Benchmark", col_span=6
        )
        row = _write_table_header(
            ws3, row, ["Period", "Your Rate", "Industry Benchmark", "Variance", "", ""]
        )
        for period in ("90d", "180d", "1y"):
            comp = ret_comp.get(period, {})
            actual = comp.get("actual")
            bench = comp.get("benchmark")
            variance = comp.get("variance")
            v_font = f_green if (variance and variance >= 0) else f_red
            row = _write_table_row(
                ws3,
                row,
                [
                    period.replace("d", " Day").replace("1y", "1 Year"),
                    f"{actual:.0f}%" if actual is not None else "N/A",
                    f"{bench:.0f}%" if bench is not None else "N/A",
                    f"{variance:+.1f}%" if variance is not None else "N/A",
                    "",
                    "",
                ],
                fonts=[f_body_bold, f_body, f_body, v_font, f_body, f_body],
            )

    # ── SHEET 4: Cost Intelligence ───────────────────────────────────────
    ws4 = wb.create_sheet("Cost Intelligence")
    _set_col_widths(ws4, [3, 22, 14, 16, 14, 14, 14, 14, 3])

    row = 2
    row = _write_section_header(ws4, row, "Cost per Quality Hire Analysis", col_span=7)
    row = _write_table_header(
        ws4,
        row,
        [
            "Source",
            "Total Cost",
            "CPH",
            "CPQH",
            "Quality %",
            "Value Ratio",
            "Cost Effective",
        ],
    )

    cost_sources = cost_analysis.get("sources", {})
    for source_name in sorted(cost_sources.keys()):
        cs = cost_sources[source_name]
        cpqh = cs.get("cost_per_quality_hire")
        vr = cs.get("value_ratio")
        effective = cs.get("is_cost_effective", False)
        eff_font = f_green if effective else f_red

        row = _write_table_row(
            ws4,
            row,
            [
                source_name,
                f"${cs.get('total_cost') or 0:,.0f}",
                f"${cs.get('cost_per_hire') or 0:,.0f}",
                f"${cpqh:,.0f}" if cpqh is not None else "N/A",
                f"{cs.get('quality_pct') or 0:.0f}%",
                f"{vr:.2f}x" if vr is not None else "N/A",
                "Yes" if effective else "No",
            ],
            fonts=[f_body_bold, f_body, f_body, f_body_bold, f_body, f_body, eff_font],
        )

    # ── SHEET 5: Recommendations ─────────────────────────────────────────
    ws5 = wb.create_sheet("Recommendations")
    _set_col_widths(ws5, [3, 6, 30, 50, 12, 12, 3])

    row = 2
    row = _write_section_header(ws5, row, "Action Plan", col_span=5)
    row = _write_table_header(
        ws5, row, ["#", "Recommendation", "Detail", "Impact", "Priority"]
    )

    for i, rec in enumerate(recommendations[:10], 1):
        priority = rec.get("priority", "medium")
        p_font = (
            f_red
            if priority in ("critical", "high")
            else (f_amber if priority == "medium" else f_body)
        )
        row = _write_table_row(
            ws5,
            row,
            [
                i,
                rec.get("title") or "",
                rec.get("description") or "",
                f"{rec.get('impact_score') or 0}/100",
                priority.upper(),
            ],
            fonts=[f_body_bold, f_body_bold, f_body, f_body, p_font],
            aligns=[align_center, align_left, align_left, align_center, align_center],
        )

    # ── SHEET 6: Predictions ──────────────────────────────────────────────
    ws6 = wb.create_sheet("Predictions")
    _set_col_widths(ws6, [3, 22, 14, 14, 14, 14, 14, 14, 3])

    row = 2
    row = _write_section_header(
        ws6, row, "Next Quarter Source Performance Predictions", col_span=7
    )
    row = _write_table_header(
        ws6,
        row,
        [
            "Source",
            "Current QoH",
            "Projected QoH",
            "Proj. Grade",
            "Trend",
            "Confidence",
            "Range",
        ],
    )

    for source_name in sorted(predictions.keys()):
        pred = predictions[source_name]
        trend = pred.get("trend", "stable")
        trend_font = (
            f_green
            if trend == "improving"
            else (f_red if trend == "declining" else f_body)
        )

        row = _write_table_row(
            ws6,
            row,
            [
                source_name,
                f"{pred.get('current_qoh') or 0:.1f}",
                f"{pred.get('projected_qoh') or 0:.1f}",
                pred.get("projected_grade", "C"),
                trend.capitalize(),
                pred.get("confidence", "low").capitalize(),
                f"{pred.get('range_low') or 0:.0f} - {pred.get('range_high') or 0:.0f}",
            ],
            fonts=[
                f_body_bold,
                f_body,
                f_body_bold,
                f_body,
                trend_font,
                f_body,
                f_body,
            ],
        )

    row += 2
    ws6.cell(
        row=row,
        column=B,
        value=f"Methodology: {report.get('predictions', {}).get('methodology', 'Weighted trend projection')}",
    ).font = f_footnote
    row += 1
    ws6.cell(
        row=row,
        column=B,
        value="Predictions are estimates based on historical patterns. Actual results may vary.",
    ).font = f_footnote

    # Save
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# =============================================================================
# 11. GENERATE SIGNAL PPT
# =============================================================================


def generate_signal_ppt(report: Dict[str, Any], client_name: str = "Client") -> bytes:
    """Generate 6-slide PPT report with HireSignal analysis.

    Slides:
        1. Title
        2. QoH Scorecard & Overview
        3. Source Effectiveness Matrix
        4. Retention & Cost Intelligence
        5. Recommendations
        6. Predictions

    Uses Joveo branding: Port Gore #202058, Blue Violet #5A54BD, Downy #6BB3CD.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_SHAPE
    except ImportError:
        logger.error("python-pptx not available for PPT generation")
        return b""

    # Joveo brand colors
    NAVY = RGBColor(0x20, 0x20, 0x58)
    BLUE = RGBColor(0x5A, 0x54, 0xBD)
    TEAL = RGBColor(0x6B, 0xB3, 0xCD)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    OFF_WHITE = RGBColor(0xFF, 0xFD, 0xF9)
    DARK_TEXT = RGBColor(0x20, 0x20, 0x58)
    MUTED_TEXT = RGBColor(0x59, 0x67, 0x80)
    GREEN = RGBColor(0x33, 0x87, 0x21)
    AMBER = RGBColor(0xCE, 0x90, 0x47)
    RED_ACCENT = RGBColor(0xB5, 0x66, 0x9C)
    LIGHT_BG = RGBColor(0xF5, 0xF3, 0xFF)
    LIGHT_GREEN = RGBColor(0xE6, 0xF2, 0xE0)
    LIGHT_RED = RGBColor(0xFD, 0xE8, 0xE8)
    WARM_GRAY = RGBColor(0xEB, 0xE6, 0xE0)

    FONT_TITLE = "Poppins"
    FONT_BODY = "Inter"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    source_eff = report.get("source_effectiveness", {})
    sources = source_eff.get("sources", {})
    retention = report.get("retention", {})
    cost_analysis = report.get("cost_analysis", {})
    benchmark = report.get("benchmark", {})
    recommendations = report.get("recommendations") or []
    predictions_data = report.get("predictions", {}).get("predictions", {})
    overall_qoh = source_eff.get("overall_qoh", 50.0)
    overall_grade = source_eff.get("overall_grade", "C")

    def _add_bg(slide, color=OFF_WHITE):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_shape(slide, left, top, width, height, fill_color, border_color=None):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1)
        else:
            shape.line.fill.background()
        return shape

    def _add_text_box(
        slide,
        left,
        top,
        width,
        height,
        text,
        font_name=FONT_BODY,
        size=12,
        color=DARK_TEXT,
        bold=False,
        align=PP_ALIGN.LEFT,
    ):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(text)
        p.font.name = font_name
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = align
        return txBox

    def _footer(slide):
        _add_shape(slide, Inches(0), Inches(7.0), prs.slide_width, Inches(0.5), NAVY)
        _add_text_box(
            slide,
            Inches(0.5),
            Inches(7.05),
            Inches(5),
            Inches(0.4),
            "Powered by HireSignal | Joveo Nova AI",
            FONT_BODY,
            10,
            WHITE,
            align=PP_ALIGN.LEFT,
        )

    # ── SLIDE 1: Title ────────────────────────────────────────────────
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide1, NAVY)
    _add_shape(slide1, Inches(0), Inches(0), prs.slide_width, Inches(0.08), TEAL)

    _add_text_box(
        slide1,
        Inches(1.5),
        Inches(2.0),
        Inches(10),
        Inches(1.2),
        "HireSignal",
        FONT_TITLE,
        44,
        WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _add_text_box(
        slide1,
        Inches(1.5),
        Inches(3.2),
        Inches(10),
        Inches(0.6),
        "Quality of Hire Signal Tracker",
        FONT_TITLE,
        20,
        TEAL,
        align=PP_ALIGN.CENTER,
    )
    _add_text_box(
        slide1,
        Inches(1.5),
        Inches(4.2),
        Inches(10),
        Inches(0.5),
        f"{client_name} | {source_eff.get('total_hires') or 0} Hires Analyzed",
        FONT_BODY,
        16,
        TEAL,
        align=PP_ALIGN.CENTER,
    )
    _add_text_box(
        slide1,
        Inches(1.5),
        Inches(5.0),
        Inches(10),
        Inches(0.5),
        f"Generated {datetime.datetime.now().strftime('%B %d, %Y')}",
        FONT_BODY,
        12,
        MUTED_TEXT,
        align=PP_ALIGN.CENTER,
    )

    _add_shape(slide1, Inches(0), Inches(7.0), prs.slide_width, Inches(0.5), BLUE)
    _add_text_box(
        slide1,
        Inches(0.5),
        Inches(7.05),
        Inches(4),
        Inches(0.4),
        "Powered by HireSignal | Joveo",
        FONT_BODY,
        10,
        WHITE,
        align=PP_ALIGN.LEFT,
    )

    # ── SLIDE 2: QoH Scorecard ────────────────────────────────────────
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide2, OFF_WHITE)

    _add_text_box(
        slide2,
        Inches(0.8),
        Inches(0.4),
        Inches(8),
        Inches(0.7),
        "Quality of Hire Overview",
        FONT_TITLE,
        28,
        NAVY,
        bold=True,
    )
    _add_shape(slide2, Inches(0.8), Inches(1.0), Inches(2), Inches(0.04), TEAL)

    # Grade card
    grade_color = (
        GREEN
        if overall_grade in ("A", "B")
        else (AMBER if overall_grade == "C" else RED_ACCENT)
    )
    grade_bg = (
        LIGHT_GREEN
        if overall_grade in ("A", "B")
        else (RGBColor(0xFE, 0xF3, 0xC7) if overall_grade == "C" else LIGHT_RED)
    )

    _add_shape(
        slide2, Inches(0.8), Inches(1.5), Inches(2.5), Inches(2.5), grade_bg, WARM_GRAY
    )
    _add_text_box(
        slide2,
        Inches(0.8),
        Inches(1.6),
        Inches(2.5),
        Inches(0.4),
        "OVERALL QoH GRADE",
        FONT_BODY,
        10,
        MUTED_TEXT,
        align=PP_ALIGN.CENTER,
    )
    _add_text_box(
        slide2,
        Inches(0.8),
        Inches(2.0),
        Inches(2.5),
        Inches(1.2),
        overall_grade,
        FONT_TITLE,
        64,
        grade_color,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _add_text_box(
        slide2,
        Inches(0.8),
        Inches(3.2),
        Inches(2.5),
        Inches(0.4),
        f"Score: {overall_qoh:.0f}/100",
        FONT_BODY,
        14,
        DARK_TEXT,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # KPI cards
    overall_cost = cost_analysis.get("overall", {})
    qoh_comp = benchmark.get("qoh_comparison", {})
    kpi_cards = [
        ("Total Hires", str(source_eff.get("total_hires") or 0)),
        ("Sources", str(source_eff.get("source_count") or 0)),
        ("Quality Hire %", f"{overall_cost.get('quality_pct') or 0:.0f}%"),
        ("CPH", f"${overall_cost.get('cost_per_hire') or 0:,.0f}"),
        (
            "CPQH",
            (
                f"${overall_cost.get('cost_per_quality_hire') or 0:,.0f}"
                if overall_cost.get("cost_per_quality_hire")
                else "N/A"
            ),
        ),
        ("vs Industry", f"{qoh_comp.get('vs_avg') or 0:+.0f} pts"),
    ]

    for i, (label, value) in enumerate(kpi_cards):
        col = i % 3
        r = i // 3
        left = Inches(3.8 + col * 3.2)
        top = Inches(1.5 + r * 1.8)
        _add_shape(slide2, left, top, Inches(2.8), Inches(1.5), WHITE, WARM_GRAY)
        _add_text_box(
            slide2,
            left + Inches(0.2),
            top + Inches(0.2),
            Inches(2.4),
            Inches(0.3),
            label,
            FONT_BODY,
            10,
            MUTED_TEXT,
        )
        _add_text_box(
            slide2,
            left + Inches(0.2),
            top + Inches(0.6),
            Inches(2.4),
            Inches(0.6),
            value,
            FONT_TITLE,
            22,
            NAVY,
            bold=True,
        )

    _add_text_box(
        slide2,
        Inches(0.8),
        Inches(5.2),
        Inches(11.5),
        Inches(0.5),
        benchmark.get("tier_description") or "",
        FONT_BODY,
        11,
        MUTED_TEXT,
    )
    _footer(slide2)

    # ── SLIDE 3: Source Effectiveness ─────────────────────────────────
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide3, OFF_WHITE)

    _add_text_box(
        slide3,
        Inches(0.8),
        Inches(0.4),
        Inches(8),
        Inches(0.7),
        "Source Effectiveness Matrix",
        FONT_TITLE,
        28,
        NAVY,
        bold=True,
    )
    _add_shape(slide3, Inches(0.8), Inches(1.0), Inches(2), Inches(0.04), TEAL)

    # Table
    sorted_sources = sorted(
        sources.values(), key=lambda x: x.get("avg_qoh_score") or 0, reverse=True
    )
    table_headers = [
        "Source",
        "Hires",
        "QoH",
        "Grade",
        "CPQH",
        "Retention",
        "Quality %",
    ]
    num_rows = min(len(sorted_sources) + 1, 10)
    num_cols = len(table_headers)

    table_shape = slide3.shapes.add_table(
        num_rows,
        num_cols,
        Inches(0.5),
        Inches(1.4),
        Inches(12.3),
        Inches(0.5 * num_rows),
    )
    table = table_shape.table

    for j, header in enumerate(table_headers):
        cell = table.cell(0, j)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.name = FONT_BODY
            paragraph.font.size = Pt(9)
            paragraph.font.color.rgb = WHITE
            paragraph.font.bold = True
            paragraph.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY

    for i, src in enumerate(sorted_sources[: num_rows - 1]):
        grade = src.get("qoh_grade", "C")
        cpqh = (
            f"${src['cost_per_quality_hire']:,.0f}"
            if src.get("cost_per_quality_hire")
            else "N/A"
        )
        ret = (
            f"{src['retention_rate']:.0f}%"
            if src.get("retention_rate") is not None
            else "N/A"
        )

        row_data = [
            src["source"],
            str(src["total_hires"]),
            f"{src['avg_qoh_score']:.0f}",
            grade,
            cpqh,
            ret,
            f"{src['quality_hire_pct']:.0f}%",
        ]

        for j, val in enumerate(row_data):
            cell = table.cell(i + 1, j)
            cell.text = val
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = FONT_BODY
                paragraph.font.size = Pt(9)
                paragraph.font.color.rgb = DARK_TEXT
                paragraph.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT

            if j == 3:  # Grade column
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = (
                        GREEN
                        if grade in ("A", "B")
                        else (AMBER if grade == "C" else RED_ACCENT)
                    )

        if i % 2 == 0:
            for j in range(num_cols):
                table.cell(i + 1, j).fill.solid()
                table.cell(i + 1, j).fill.fore_color.rgb = LIGHT_BG

    _footer(slide3)

    # ── SLIDE 4: Retention & Cost ─────────────────────────────────────
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide4, OFF_WHITE)

    _add_text_box(
        slide4,
        Inches(0.8),
        Inches(0.4),
        Inches(8),
        Inches(0.7),
        "Retention & Cost Intelligence",
        FONT_TITLE,
        28,
        NAVY,
        bold=True,
    )
    _add_shape(slide4, Inches(0.8), Inches(1.0), Inches(2), Inches(0.04), TEAL)

    # Retention cards - top sources
    _add_text_box(
        slide4,
        Inches(0.8),
        Inches(1.3),
        Inches(5),
        Inches(0.3),
        "RETENTION BY SOURCE",
        FONT_BODY,
        10,
        MUTED_TEXT,
        bold=True,
    )

    max_cards = min(len(sorted_sources), 4)
    for i, src in enumerate(sorted_sources[:max_cards]):
        left = Inches(0.5 + i * 3.1)
        top = Inches(1.8)
        grade = src.get("qoh_grade", "C")
        card_border = (
            GREEN if grade in ("A", "B") else (AMBER if grade == "C" else RED_ACCENT)
        )
        _add_shape(slide4, left, top, Inches(2.8), Inches(2.0), WHITE, card_border)

        _add_text_box(
            slide4,
            left + Inches(0.15),
            top + Inches(0.1),
            Inches(2.5),
            Inches(0.3),
            src["source"],
            FONT_TITLE,
            12,
            NAVY,
            bold=True,
        )
        ret = src.get("retention_rate")
        ret_text = f"{ret:.0f}%" if ret is not None else "N/A"
        _add_text_box(
            slide4,
            left + Inches(0.15),
            top + Inches(0.5),
            Inches(2.5),
            Inches(0.5),
            ret_text,
            FONT_TITLE,
            28,
            card_border,
            bold=True,
        )
        _add_text_box(
            slide4,
            left + Inches(0.15),
            top + Inches(1.1),
            Inches(2.5),
            Inches(0.3),
            f"QoH: {src['avg_qoh_score']:.0f} | Hires: {src['total_hires']}",
            FONT_BODY,
            9,
            MUTED_TEXT,
        )
        cpqh = src.get("cost_per_quality_hire")
        _add_text_box(
            slide4,
            left + Inches(0.15),
            top + Inches(1.5),
            Inches(2.5),
            Inches(0.3),
            f"CPQH: ${cpqh:,.0f}" if cpqh else "CPQH: N/A",
            FONT_BODY,
            9,
            MUTED_TEXT,
        )

    # Cost comparison at bottom
    _add_text_box(
        slide4,
        Inches(0.8),
        Inches(4.2),
        Inches(5),
        Inches(0.3),
        "COST PER HIRE vs COST PER QUALITY HIRE",
        FONT_BODY,
        10,
        MUTED_TEXT,
        bold=True,
    )

    cost_sources_list = sorted(
        cost_analysis.get("sources", {}).values(),
        key=lambda x: x.get("cost_per_hire") or 0,
    )
    for i, cs in enumerate(cost_sources_list[:4]):
        left = Inches(0.5 + i * 3.1)
        top = Inches(4.7)
        _add_shape(slide4, left, top, Inches(2.8), Inches(1.8), WHITE, WARM_GRAY)
        _add_text_box(
            slide4,
            left + Inches(0.15),
            top + Inches(0.1),
            Inches(2.5),
            Inches(0.25),
            cs["source"],
            FONT_BODY,
            10,
            NAVY,
            bold=True,
        )
        _add_text_box(
            slide4,
            left + Inches(0.15),
            top + Inches(0.4),
            Inches(2.5),
            Inches(0.3),
            f"CPH: ${cs.get('cost_per_hire') or 0:,.0f}",
            FONT_BODY,
            12,
            DARK_TEXT,
            bold=True,
        )
        cpqh = cs.get("cost_per_quality_hire")
        _add_text_box(
            slide4,
            left + Inches(0.15),
            top + Inches(0.8),
            Inches(2.5),
            Inches(0.3),
            f"CPQH: ${cpqh:,.0f}" if cpqh else "CPQH: N/A",
            FONT_BODY,
            12,
            BLUE,
            bold=True,
        )
        vr = cs.get("value_ratio")
        effective = cs.get("is_cost_effective", False)
        _add_text_box(
            slide4,
            left + Inches(0.15),
            top + Inches(1.2),
            Inches(2.5),
            Inches(0.3),
            (
                f"Ratio: {vr:.1f}x {'(Efficient)' if effective else '(Needs work)'}"
                if vr
                else ""
            ),
            FONT_BODY,
            9,
            GREEN if effective else RED_ACCENT,
        )

    _footer(slide4)

    # ── SLIDE 5: Recommendations ──────────────────────────────────────
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide5, OFF_WHITE)

    _add_text_box(
        slide5,
        Inches(0.8),
        Inches(0.4),
        Inches(8),
        Inches(0.7),
        "Action Plan",
        FONT_TITLE,
        28,
        NAVY,
        bold=True,
    )
    _add_shape(slide5, Inches(0.8), Inches(1.0), Inches(2), Inches(0.04), TEAL)

    max_recs = min(len(recommendations), 6)
    for i, rec in enumerate(recommendations[:max_recs]):
        priority = rec.get("priority", "medium")
        top = Inches(1.4 + i * 0.9)

        if priority in ("critical", "high"):
            action_color = RED_ACCENT
            bg_color = LIGHT_RED
        elif priority == "medium":
            action_color = AMBER
            bg_color = RGBColor(0xFE, 0xF3, 0xC7)
        else:
            action_color = GREEN
            bg_color = LIGHT_GREEN

        _add_shape(
            slide5, Inches(0.5), top, Inches(12.3), Inches(0.75), bg_color, action_color
        )

        # Impact badge
        impact = rec.get("impact_score") or 0
        _add_shape(
            slide5,
            Inches(0.7),
            top + Inches(0.12),
            Inches(1.2),
            Inches(0.45),
            action_color,
        )
        _add_text_box(
            slide5,
            Inches(0.7),
            top + Inches(0.15),
            Inches(1.2),
            Inches(0.4),
            f"Impact: {impact}",
            FONT_BODY,
            9,
            WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )

        # Title
        _add_text_box(
            slide5,
            Inches(2.1),
            top + Inches(0.08),
            Inches(10),
            Inches(0.35),
            rec.get("title") or "",
            FONT_TITLE,
            11,
            NAVY,
            bold=True,
        )
        # Description (truncated)
        desc = rec.get("description") or ""
        if len(desc) > 120:
            desc = desc[:117] + "..."
        _add_text_box(
            slide5,
            Inches(2.1),
            top + Inches(0.38),
            Inches(10),
            Inches(0.35),
            desc,
            FONT_BODY,
            8,
            MUTED_TEXT,
        )

    _footer(slide5)

    # ── SLIDE 6: Predictions ──────────────────────────────────────────
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide6, OFF_WHITE)

    _add_text_box(
        slide6,
        Inches(0.8),
        Inches(0.4),
        Inches(8),
        Inches(0.7),
        "Next Quarter Predictions",
        FONT_TITLE,
        28,
        NAVY,
        bold=True,
    )
    _add_shape(slide6, Inches(0.8), Inches(1.0), Inches(2), Inches(0.04), TEAL)

    sorted_preds = sorted(
        predictions_data.values(),
        key=lambda x: x.get("projected_qoh") or 0,
        reverse=True,
    )
    max_pred_cards = min(len(sorted_preds), 6)

    for i, pred in enumerate(sorted_preds[:max_pred_cards]):
        col = i % 3
        r = i // 3
        left = Inches(0.5 + col * 4.2)
        top = Inches(1.5 + r * 2.8)

        trend = pred.get("trend", "stable")
        card_border = (
            GREEN
            if trend == "improving"
            else (RED_ACCENT if trend == "declining" else WARM_GRAY)
        )
        _add_shape(slide6, left, top, Inches(3.8), Inches(2.4), WHITE, card_border)

        _add_text_box(
            slide6,
            left + Inches(0.2),
            top + Inches(0.15),
            Inches(2.8),
            Inches(0.35),
            pred.get("source") or "",
            FONT_TITLE,
            14,
            NAVY,
            bold=True,
        )

        trend_symbol = (
            "Improving"
            if trend == "improving"
            else ("Declining" if trend == "declining" else "Stable")
        )
        trend_color = (
            GREEN
            if trend == "improving"
            else (RED_ACCENT if trend == "declining" else MUTED_TEXT)
        )
        _add_text_box(
            slide6,
            left + Inches(3.0),
            top + Inches(0.15),
            Inches(0.6),
            Inches(0.35),
            trend_symbol[:3],
            FONT_BODY,
            10,
            trend_color,
            bold=True,
            align=PP_ALIGN.CENTER,
        )

        pred_lines = [
            f"Current QoH: {pred.get('current_qoh') or 0:.0f}  ->  Projected: {pred.get('projected_qoh') or 0:.0f}",
            f"Grade: {pred.get('projected_grade', 'C')}  |  Confidence: {pred.get('confidence', 'low').capitalize()}",
            f"Range: {pred.get('range_low') or 0:.0f} - {pred.get('range_high') or 0:.0f}",
            f"Projected Retention: {pred.get('projected_retention') or 0:.0f}%",
        ]
        for m_idx, m_text in enumerate(pred_lines):
            _add_text_box(
                slide6,
                left + Inches(0.2),
                top + Inches(0.6 + m_idx * 0.4),
                Inches(3.4),
                Inches(0.35),
                m_text,
                FONT_BODY,
                9,
                MUTED_TEXT if m_idx < 3 else DARK_TEXT,
                bold=(m_idx == 3),
            )

    _footer(slide6)

    # Save
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# =============================================================================
# 12. ORCHESTRATOR -- single entry point
# =============================================================================


def run_full_signal_analysis(
    file_bytes: Optional[bytes] = None,
    filename: str = "",
    industry: str = "general_entry_level",
    client_name: str = "Client",
    manual_data: Optional[List[Dict]] = None,
) -> Dict[str, Any]:
    """Full HireSignal analysis pipeline.

    Accepts either file upload (file_bytes + filename) or manual_data (list of hire dicts).

    Pipeline:
        1. Parse hiring data (or use manual_data)
        2. Calculate QoH scores for each hire
        3. Analyze source effectiveness
        4. Calculate retention rates
        5. Calculate CPQH
        6. Benchmark against industry
        7. Generate recommendations
        8. Predict source performance

    Returns complete report_data dict ready for Excel/PPT generation.
    Thread-safe, never raises.
    """
    try:
        # 1. Parse or accept data
        if manual_data:
            hires = manual_data
        elif file_bytes:
            hires = parse_hiring_data(file_bytes, filename)
        else:
            return {
                "error": "No hiring data provided. Upload a file or provide manual data.",
                "success": False,
            }

        if not hires:
            return {
                "error": "Could not parse hiring data. Please check file format and column headers.",
                "success": False,
            }

        # 2. Calculate QoH score for each hire (guard against non-dict entries)
        for i, h in enumerate(hires):
            if not isinstance(h, dict):
                logger.warning(
                    "hires[%d] is %s, not dict -- skipping", i, type(h).__name__
                )
                continue
            h["qoh_score"] = calculate_qoh_score(h)

        # Use ThreadPoolExecutor for concurrent analysis
        results: Dict[str, Any] = {}
        with ThreadPoolExecutor(max_workers=4) as executor:
            futures = {
                executor.submit(
                    analyze_source_effectiveness, hires
                ): "source_effectiveness",
                executor.submit(calculate_retention_rates, hires): "retention",
                executor.submit(
                    calculate_cost_per_quality_hire, hires
                ): "cost_analysis",
            }
            for future in as_completed(futures):
                key = futures[future]
                try:
                    results[key] = future.result()
                except Exception as exc:
                    logger.error("HireSignal %s analysis failed: %s", key, exc)
                    results[key] = {}

        # 6. Benchmark (needs source_effectiveness results)
        try:
            results["benchmark"] = benchmark_against_industry(
                {**results, "retention": results.get("retention", {})},
                industry,
            )
        except Exception as exc:
            logger.error("Benchmark comparison failed: %s", exc)
            results["benchmark"] = {}

        # 7. Recommendations (needs all prior results)
        try:
            results["recommendations"] = generate_recommendations(results)
        except Exception as exc:
            logger.error("Recommendation generation failed: %s", exc)
            results["recommendations"] = []

        # 8. Predictions
        try:
            results["predictions"] = predict_source_performance(results)
        except Exception as exc:
            logger.error("Prediction generation failed: %s", exc)
            results["predictions"] = {"predictions": {}, "confidence": "low"}

        return {
            "success": True,
            "client_name": client_name,
            "industry": industry,
            "industry_label": INDUSTRY_LABEL_MAP.get(
                industry, industry.replace("_", " ").title()
            ),
            "total_hires": len(hires),
            "hires": hires[:500],  # Cap for JSON response size
            "source_effectiveness": results.get("source_effectiveness", {}),
            "retention": results.get("retention", {}),
            "cost_analysis": results.get("cost_analysis", {}),
            "benchmark": results.get("benchmark", {}),
            "recommendations": results.get("recommendations") or [],
            "predictions": results.get("predictions", {}),
            "generated_at": datetime.datetime.now().isoformat(),
        }

    except Exception as exc:
        logger.exception("HireSignal analysis failed: %s", exc)
        return {"error": f"Analysis failed: {str(exc)}", "success": False}
