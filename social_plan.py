"""
social_plan.py -- Social & Search Media Plan Generator

Generates recruitment media plans for social channels (Facebook, Instagram,
LinkedIn, TikTok, Twitter/X, Snapchat, Reddit, YouTube, Pinterest) and search
channels (Google Ads, Bing Ads, Indeed Sponsored). Uses embedded KB data +
trend_engine benchmarks + collar_intelligence for collar-aware routing.

Architecture mirrors quick_plan.py: lazy imports, try/except, ThreadPoolExecutor
for concurrent platform analysis, zero external API calls at runtime.

Thread-safe, never crashes (all exceptions caught and degraded gracefully).
"""

from __future__ import annotations

import io
import json
import logging
import math
import os
import re
import threading
from concurrent.futures import ThreadPoolExecutor, as_completed
from datetime import datetime, timedelta
from typing import Any, Dict, List, Optional, Tuple

logger = logging.getLogger(__name__)

_LOCK = threading.Lock()
_BASE_DIR = os.path.dirname(os.path.abspath(__file__))

# ── Lazy imports ──

try:
    import trend_engine as _trend_engine

    _HAS_TREND_ENGINE = True
except ImportError:
    _trend_engine = None
    _HAS_TREND_ENGINE = False

try:
    import collar_intelligence as _collar_intel

    _HAS_COLLAR_INTEL = True
except ImportError:
    _collar_intel = None
    _HAS_COLLAR_INTEL = False

try:
    import research as _research

    _HAS_RESEARCH = True
except ImportError:
    _research = None
    _HAS_RESEARCH = False

try:
    from shared_utils import (
        INDUSTRY_LABEL_MAP,
        parse_budget,
        standardize_location,
        parse_budget_display,
    )
except ImportError:
    INDUSTRY_LABEL_MAP = {}

    def parse_budget(v, *, default=100_000.0):
        try:
            return float(v)
        except (TypeError, ValueError):
            return default

    def standardize_location(s):
        return s

    def parse_budget_display(v):
        try:
            return float(v)
        except (TypeError, ValueError):
            return None


try:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter
    from openpyxl.chart import PieChart, BarChart, Reference
    from openpyxl.chart.label import DataLabelList
    from openpyxl.chart.series import DataPoint

    _HAS_OPENPYXL = True
except ImportError:
    _HAS_OPENPYXL = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE

    _HAS_PPTX = True
except ImportError:
    _HAS_PPTX = False


# ═══════════════════════════════════════════════════════════════════════════════
# PLATFORM DEFINITIONS
# ═══════════════════════════════════════════════════════════════════════════════

SOCIAL_PLATFORMS: Dict[str, Dict[str, Any]] = {
    "facebook": {
        "name": "Facebook",
        "icon": "fb",
        "color": "#1877F2",
        "monthly_active_users": "3.07B",
        "demographics": "25-54 primary, broad reach",
        "best_for_collar": ["blue_collar", "grey_collar", "pink_collar"],
        "ad_formats": [
            "Image",
            "Video",
            "Carousel",
            "Stories",
            "Reels",
            "Lead Gen Forms",
            "Collection",
        ],
        "min_daily_budget": 5.0,
        "trend_engine_key": "meta_facebook",
        "strengths": [
            "Massive reach",
            "Granular targeting",
            "Lead Gen Forms",
            "Lookalike audiences",
            "Low CPC for blue collar",
        ],
        "weaknesses": ["Declining organic reach", "Ad fatigue", "Privacy restrictions"],
        "best_times": {"weekdays": "9-11 AM, 1-3 PM", "weekends": "10 AM-12 PM"},
        "image_specs": {
            "feed": "1080x1080 or 1200x628",
            "stories": "1080x1920",
            "video": "1080x1080 or 16:9",
        },
        "copy_limits": {"primary": 125, "headline": 40, "description": 30},
        "avg_engagement_rate": 0.06,
    },
    "instagram": {
        "name": "Instagram",
        "icon": "ig",
        "color": "#E4405F",
        "monthly_active_users": "2.4B",
        "demographics": "18-34 primary, visual-first",
        "best_for_collar": ["white_collar", "pink_collar", "grey_collar"],
        "ad_formats": [
            "Stories",
            "Reels",
            "Feed Image",
            "Carousel",
            "Explore",
            "Shopping",
        ],
        "min_daily_budget": 5.0,
        "trend_engine_key": "meta_instagram",
        "strengths": [
            "High engagement",
            "Visual storytelling",
            "Reels momentum",
            "Younger demographics",
        ],
        "weaknesses": [
            "Higher CPC than Facebook",
            "Limited link options",
            "Not ideal for blue collar",
        ],
        "best_times": {"weekdays": "11 AM-1 PM, 7-9 PM", "weekends": "10 AM-2 PM"},
        "image_specs": {
            "feed": "1080x1080",
            "stories": "1080x1920",
            "reels": "1080x1920",
        },
        "copy_limits": {"primary": 125, "headline": 40, "description": 30},
        "avg_engagement_rate": 0.08,
    },
    "linkedin": {
        "name": "LinkedIn",
        "icon": "li",
        "color": "#0A66C2",
        "monthly_active_users": "310M MAU (1.3B members)",
        "demographics": "25-54, professional, college-educated",
        "best_for_collar": ["white_collar"],
        "ad_formats": [
            "Sponsored Content",
            "Message Ads",
            "Dynamic Ads",
            "Text Ads",
            "Video Ads",
            "Document Ads",
            "Lead Gen Forms",
        ],
        "min_daily_budget": 10.0,
        "trend_engine_key": "linkedin",
        "strengths": [
            "Highest quality professional candidates",
            "Job title targeting",
            "Company targeting",
            "InMail",
            "89% of recruiters use it",
        ],
        "weaknesses": [
            "Highest CPC",
            "Limited blue collar reach",
            "Ad fatigue among power users",
        ],
        "best_times": {
            "weekdays": "7-8 AM, 12 PM, 5-6 PM (Tue-Thu best)",
            "weekends": "Not recommended",
        },
        "image_specs": {
            "feed": "1200x627",
            "stories": "1080x1920",
            "carousel": "1080x1080",
        },
        "copy_limits": {"intro": 150, "headline": 70, "description": 100},
        "avg_engagement_rate": 0.04,
    },
    "tiktok": {
        "name": "TikTok",
        "icon": "tt",
        "color": "#010101",
        "monthly_active_users": "1.6B",
        "demographics": "16-34 primary, Gen Z dominant",
        "best_for_collar": ["blue_collar", "pink_collar"],
        "ad_formats": [
            "In-Feed Video",
            "TopView",
            "Branded Hashtag",
            "Spark Ads",
            "Lead Gen",
        ],
        "min_daily_budget": 20.0,
        "trend_engine_key": None,
        "strengths": [
            "Gen Z reach",
            "Viral potential",
            "Authentic content",
            "Low CPM",
            "High engagement",
        ],
        "weaknesses": [
            "Limited professional targeting",
            "Short content lifespan",
            "Brand safety concerns",
        ],
        "best_times": {
            "weekdays": "7-9 AM, 12-3 PM, 7-11 PM",
            "weekends": "9 AM-12 PM",
        },
        "image_specs": {
            "video": "1080x1920 (9:16)",
            "min_duration": "5s",
            "max_duration": "60s recommended",
        },
        "copy_limits": {"description": 100, "display_name": 40},
        "avg_engagement_rate": 0.12,
        "benchmarks_2025": {
            "avg_cpc": 0.50,
            "avg_cpm": 6.00,
            "avg_ctr": 0.018,
            "avg_cpa": 8.00,
        },
    },
    "twitter_x": {
        "name": "Twitter / X",
        "icon": "x",
        "color": "#000000",
        "monthly_active_users": "600M",
        "demographics": "25-49, news/tech-savvy",
        "best_for_collar": ["white_collar"],
        "ad_formats": [
            "Promoted Posts",
            "Video Ads",
            "Carousel",
            "Follower Ads",
            "Trend Takeover",
        ],
        "min_daily_budget": 1.0,
        "trend_engine_key": None,
        "strengths": [
            "Real-time engagement",
            "Tech audience",
            "Conversation targeting",
            "Low minimum spend",
        ],
        "weaknesses": [
            "Brand safety concerns",
            "Declining ad revenue",
            "Uncertain platform future",
        ],
        "best_times": {"weekdays": "8-10 AM, 12-1 PM", "weekends": "9 AM"},
        "image_specs": {"feed": "1200x675", "video": "1280x720 min"},
        "copy_limits": {"text": 280, "headline": 70},
        "avg_engagement_rate": 0.03,
        "benchmarks_2025": {
            "avg_cpc": 0.80,
            "avg_cpm": 5.50,
            "avg_ctr": 0.012,
            "avg_cpa": 15.00,
        },
    },
    "snapchat": {
        "name": "Snapchat",
        "icon": "sc",
        "color": "#FFFC00",
        "monthly_active_users": "850M",
        "demographics": "13-34, Gen Z/young Millennial",
        "best_for_collar": ["blue_collar", "pink_collar"],
        "ad_formats": [
            "Snap Ads",
            "Story Ads",
            "Collection Ads",
            "Filters",
            "Lenses",
            "Commercials",
        ],
        "min_daily_budget": 5.0,
        "trend_engine_key": None,
        "strengths": [
            "Young audience",
            "AR features",
            "Low competition for recruitment",
            "Full-screen immersive",
        ],
        "weaknesses": [
            "Small professional audience",
            "Limited targeting",
            "High skip rates",
        ],
        "best_times": {"weekdays": "10 PM-1 AM", "weekends": "All day"},
        "image_specs": {"snap_ad": "1080x1920", "video": "1080x1920 (3-180s)"},
        "copy_limits": {"headline": 34, "brand_name": 25},
        "avg_engagement_rate": 0.05,
        "benchmarks_2025": {
            "avg_cpc": 0.45,
            "avg_cpm": 4.50,
            "avg_ctr": 0.015,
            "avg_cpa": 9.00,
        },
    },
    "reddit": {
        "name": "Reddit",
        "icon": "rd",
        "color": "#FF4500",
        "monthly_active_users": "1.1B",
        "demographics": "18-44, tech-savvy, niche communities",
        "best_for_collar": ["white_collar"],
        "ad_formats": [
            "Promoted Posts",
            "Video Ads",
            "Carousel",
            "Conversation Ads",
            "Free-Form Ads",
        ],
        "min_daily_budget": 5.0,
        "trend_engine_key": None,
        "strengths": [
            "Highly engaged communities",
            "Interest-based targeting",
            "Authenticity valued",
            "Low CPM",
        ],
        "weaknesses": [
            "Anti-advertising culture",
            "Limited scale",
            "Requires authentic tone",
        ],
        "best_times": {"weekdays": "6-8 AM, 12-2 PM", "weekends": "8-11 AM"},
        "image_specs": {"feed": "1200x628", "video": "up to 15 min"},
        "copy_limits": {"title": 300, "body": 40000},
        "avg_engagement_rate": 0.07,
        "benchmarks_2025": {
            "avg_cpc": 0.65,
            "avg_cpm": 3.50,
            "avg_ctr": 0.010,
            "avg_cpa": 12.00,
        },
    },
    "youtube": {
        "name": "YouTube",
        "icon": "yt",
        "color": "#FF0000",
        "monthly_active_users": "2.5B",
        "demographics": "18-65+, broadest reach",
        "best_for_collar": ["white_collar", "blue_collar", "grey_collar"],
        "ad_formats": [
            "Skippable In-Stream",
            "Non-Skippable In-Stream",
            "Bumper (6s)",
            "Discovery",
            "Shorts",
        ],
        "min_daily_budget": 10.0,
        "trend_engine_key": None,
        "strengths": [
            "Massive reach",
            "Video storytelling",
            "Google targeting data",
            "Shorts growing",
            "Employer brand building",
        ],
        "weaknesses": ["High production cost", "Skip rates", "Longer time to results"],
        "best_times": {"weekdays": "12-4 PM, 8-11 PM", "weekends": "9 AM-12 PM"},
        "image_specs": {
            "video": "1920x1080 (16:9)",
            "shorts": "1080x1920 (9:16)",
            "thumbnail": "1280x720",
        },
        "copy_limits": {"title": 100, "description": 5000},
        "avg_engagement_rate": 0.05,
        "benchmarks_2025": {
            "avg_cpc": 0.35,
            "avg_cpm": 8.00,
            "avg_ctr": 0.008,
            "avg_cpa": 18.00,
        },
    },
    "pinterest": {
        "name": "Pinterest",
        "icon": "pn",
        "color": "#E60023",
        "monthly_active_users": "537M",
        "demographics": "25-44, 60% female, lifestyle/creative",
        "best_for_collar": ["pink_collar", "white_collar"],
        "ad_formats": ["Standard Pin", "Video Pin", "Carousel", "Shopping", "Idea Ads"],
        "min_daily_budget": 5.0,
        "trend_engine_key": None,
        "strengths": [
            "High purchase intent",
            "Long content lifespan",
            "Visual discovery",
            "Low competition",
        ],
        "weaknesses": [
            "Small audience vs others",
            "Limited B2B targeting",
            "Niche use case for recruitment",
        ],
        "best_times": {"weekdays": "2-4 PM, 8-11 PM", "weekends": "8-11 AM"},
        "image_specs": {"pin": "1000x1500 (2:3)", "video": "1000x1500 or 1080x1920"},
        "copy_limits": {"title": 100, "description": 500},
        "avg_engagement_rate": 0.06,
        "benchmarks_2025": {
            "avg_cpc": 0.55,
            "avg_cpm": 5.00,
            "avg_ctr": 0.009,
            "avg_cpa": 14.00,
        },
    },
}

SEARCH_PLATFORMS: Dict[str, Dict[str, Any]] = {
    "google_ads": {
        "name": "Google Ads",
        "icon": "ga",
        "color": "#4285F4",
        "type": "search",
        "trend_engine_key": "google_search",
        "strengths": [
            "Highest intent traffic",
            "Massive reach",
            "Keyword targeting",
            "Remarketing",
        ],
        "weaknesses": ["Highest CPC", "Complex to manage", "Competitive bidding"],
        "ad_formats": [
            "Search Text Ads",
            "Responsive Search",
            "Display",
            "Video (YouTube)",
            "Performance Max",
        ],
        "best_for_collar": ["white_collar", "grey_collar", "blue_collar"],
    },
    "bing_ads": {
        "name": "Microsoft Advertising (Bing)",
        "icon": "ba",
        "color": "#00A4EF",
        "type": "search",
        "trend_engine_key": None,
        "strengths": [
            "Lower CPC than Google",
            "Older/higher income audience",
            "LinkedIn targeting integration",
        ],
        "weaknesses": ["Lower volume", "Smaller audience", "Often overlooked"],
        "ad_formats": [
            "Search Text Ads",
            "Responsive Search",
            "Audience Ads",
            "Shopping",
        ],
        "best_for_collar": ["white_collar"],
        "benchmarks_2025": {
            "avg_cpc": 1.80,
            "avg_cpm": 7.00,
            "avg_ctr": 0.032,
            "avg_cpa": 22.00,
        },
    },
    "indeed_sponsored": {
        "name": "Indeed Sponsored Jobs",
        "icon": "id",
        "color": "#2164F3",
        "type": "job_board_search",
        "trend_engine_key": "indeed",
        "strengths": [
            "Highest job seeker intent",
            "CPA model",
            "360M monthly visitors",
            "Programmatic compatible",
        ],
        "weaknesses": [
            "CPA model unpredictable",
            "Quality varies",
            "Declining organic reach",
        ],
        "ad_formats": ["Sponsored Jobs", "Indeed Resume Search", "Company Pages"],
        "best_for_collar": ["blue_collar", "grey_collar", "white_collar"],
    },
}


# ═══════════════════════════════════════════════════════════════════════════════
# COLLAR-SPECIFIC PLATFORM RANKINGS
# Score 0-100 for how well each platform serves each collar type
# ═══════════════════════════════════════════════════════════════════════════════

COLLAR_PLATFORM_SCORES: Dict[str, Dict[str, int]] = {
    "blue_collar": {
        "facebook": 95,
        "tiktok": 88,
        "instagram": 65,
        "youtube": 72,
        "snapchat": 70,
        "twitter_x": 30,
        "linkedin": 20,
        "reddit": 35,
        "pinterest": 15,
        "google_ads": 85,
        "bing_ads": 40,
        "indeed_sponsored": 92,
    },
    "white_collar": {
        "linkedin": 95,
        "google_ads": 90,
        "facebook": 60,
        "instagram": 55,
        "twitter_x": 65,
        "youtube": 68,
        "reddit": 58,
        "pinterest": 30,
        "tiktok": 35,
        "snapchat": 20,
        "bing_ads": 72,
        "indeed_sponsored": 78,
    },
    "grey_collar": {
        "facebook": 88,
        "indeed_sponsored": 90,
        "google_ads": 82,
        "linkedin": 55,
        "instagram": 60,
        "youtube": 65,
        "tiktok": 55,
        "twitter_x": 35,
        "reddit": 40,
        "snapchat": 40,
        "bing_ads": 50,
        "pinterest": 20,
    },
    "pink_collar": {
        "facebook": 92,
        "instagram": 78,
        "tiktok": 72,
        "indeed_sponsored": 85,
        "google_ads": 80,
        "snapchat": 60,
        "youtube": 58,
        "pinterest": 48,
        "linkedin": 40,
        "twitter_x": 30,
        "reddit": 25,
        "bing_ads": 45,
    },
}

# ═══════════════════════════════════════════════════════════════════════════════
# GOAL WEIGHTS -- how campaign goals shift platform priority
# ═══════════════════════════════════════════════════════════════════════════════

GOAL_PLATFORM_WEIGHTS: Dict[str, Dict[str, float]] = {
    "brand_awareness": {
        "facebook": 1.15,
        "instagram": 1.25,
        "youtube": 1.30,
        "tiktok": 1.35,
        "linkedin": 1.10,
        "twitter_x": 1.10,
        "snapchat": 1.15,
        "reddit": 0.90,
        "pinterest": 1.05,
        "google_ads": 0.80,
        "bing_ads": 0.70,
        "indeed_sponsored": 0.60,
    },
    "job_applications": {
        "indeed_sponsored": 1.40,
        "google_ads": 1.30,
        "facebook": 1.20,
        "linkedin": 1.15,
        "bing_ads": 1.10,
        "instagram": 0.95,
        "tiktok": 0.85,
        "youtube": 0.80,
        "twitter_x": 0.75,
        "reddit": 0.80,
        "snapchat": 0.70,
        "pinterest": 0.60,
    },
    "talent_pipeline": {
        "linkedin": 1.40,
        "facebook": 1.10,
        "instagram": 1.05,
        "youtube": 1.15,
        "google_ads": 1.00,
        "indeed_sponsored": 1.00,
        "twitter_x": 0.95,
        "reddit": 0.90,
        "tiktok": 0.85,
        "bing_ads": 0.85,
        "snapchat": 0.70,
        "pinterest": 0.65,
    },
    "employer_brand": {
        "linkedin": 1.35,
        "instagram": 1.30,
        "youtube": 1.35,
        "tiktok": 1.20,
        "facebook": 1.10,
        "twitter_x": 1.05,
        "pinterest": 1.00,
        "reddit": 0.95,
        "snapchat": 0.85,
        "google_ads": 0.70,
        "bing_ads": 0.60,
        "indeed_sponsored": 0.55,
    },
}

# ═══════════════════════════════════════════════════════════════════════════════
# INDUSTRY-SPECIFIC PLATFORM BOOSTS
# ═══════════════════════════════════════════════════════════════════════════════

INDUSTRY_PLATFORM_BOOSTS: Dict[str, Dict[str, float]] = {
    "tech_engineering": {
        "linkedin": 1.3,
        "twitter_x": 1.2,
        "reddit": 1.3,
        "youtube": 1.1,
    },
    "healthcare_medical": {
        "facebook": 1.2,
        "indeed_sponsored": 1.2,
        "google_ads": 1.15,
    },
    "retail_consumer": {
        "facebook": 1.2,
        "tiktok": 1.3,
        "instagram": 1.2,
        "snapchat": 1.2,
    },
    "blue_collar_trades": {"facebook": 1.25, "tiktok": 1.2, "indeed_sponsored": 1.2},
    "hospitality_travel": {"instagram": 1.3, "tiktok": 1.25, "facebook": 1.15},
    "finance_banking": {"linkedin": 1.3, "google_ads": 1.15, "twitter_x": 1.1},
    "logistics_supply_chain": {"facebook": 1.2, "indeed_sponsored": 1.25},
    "food_beverage": {
        "tiktok": 1.3,
        "instagram": 1.25,
        "facebook": 1.15,
        "snapchat": 1.15,
    },
    "construction_real_estate": {"facebook": 1.2, "indeed_sponsored": 1.2},
    "media_entertainment": {
        "instagram": 1.3,
        "tiktok": 1.3,
        "youtube": 1.25,
        "twitter_x": 1.15,
    },
    "education": {"linkedin": 1.2, "facebook": 1.15, "youtube": 1.15},
    "pharma_biotech": {"linkedin": 1.25, "google_ads": 1.15},
    "aerospace_defense": {"linkedin": 1.25, "google_ads": 1.15, "reddit": 1.1},
    "general_entry_level": {"facebook": 1.2, "tiktok": 1.25, "indeed_sponsored": 1.2},
    "automotive": {"facebook": 1.15, "youtube": 1.2, "tiktok": 1.15},
    "energy_utilities": {"linkedin": 1.15, "google_ads": 1.15, "indeed_sponsored": 1.1},
    "insurance": {"linkedin": 1.25, "google_ads": 1.15},
    "mental_health": {"facebook": 1.15, "indeed_sponsored": 1.15, "linkedin": 1.1},
    "telecommunications": {"linkedin": 1.2, "google_ads": 1.1},
}


# ═══════════════════════════════════════════════════════════════════════════════
# AUDIENCE TARGETING TEMPLATES
# ═══════════════════════════════════════════════════════════════════════════════

TARGETING_TEMPLATES: Dict[str, Dict[str, Any]] = {
    "facebook": {
        "targeting_types": [
            "Job Title",
            "Employer",
            "Industry",
            "Interests",
            "Behaviors",
            "Lookalike",
            "Custom Audience",
        ],
        "blue_collar": {
            "interests": [
                "Trade jobs",
                "Construction",
                "CDL",
                "Warehouse work",
                "Manufacturing jobs",
            ],
            "behaviors": ["Recently moved", "Job seekers", "Active job seeker"],
            "demographics": {"age": "18-55", "education": "High school+"},
            "radius_miles": 25,
        },
        "white_collar": {
            "interests": [
                "Career development",
                "Professional networking",
                "Industry news",
            ],
            "behaviors": ["Frequent travelers", "Technology early adopters"],
            "demographics": {"age": "22-55", "education": "College+"},
            "radius_miles": 50,
        },
        "grey_collar": {
            "interests": [
                "Healthcare careers",
                "Nursing",
                "Medical technology",
                "Patient care",
            ],
            "behaviors": ["Job seekers", "Healthcare workers"],
            "demographics": {"age": "22-55", "education": "Associate+"},
            "radius_miles": 30,
        },
    },
    "linkedin": {
        "targeting_types": [
            "Job Title",
            "Job Function",
            "Seniority",
            "Company Size",
            "Industry",
            "Skills",
            "Groups",
            "Company",
        ],
        "white_collar": {
            "seniority": ["Entry", "Senior", "Manager", "Director", "VP"],
            "functions": [
                "Engineering",
                "IT",
                "Finance",
                "Marketing",
                "Operations",
                "HR",
            ],
            "company_size": [
                "51-200",
                "201-500",
                "501-1000",
                "1001-5000",
                "5001-10000",
                "10000+",
            ],
            "skills_targeting": True,
        },
        "blue_collar": {
            "seniority": ["Entry", "Senior"],
            "functions": ["Operations", "Engineering", "Support"],
            "company_size": ["51-200", "201-500", "501-1000"],
            "skills_targeting": True,
        },
    },
    "tiktok": {
        "targeting_types": [
            "Interest",
            "Behavior",
            "Hashtag",
            "Creator",
            "Custom Audience",
            "Lookalike",
        ],
        "blue_collar": {
            "interests": [
                "Jobs & career",
                "Trade skills",
                "Side hustles",
                "Money tips",
            ],
            "hashtags": [
                "#hiring",
                "#jobsearch",
                "#nowhiring",
                "#worklife",
                "#dayinthelife",
            ],
            "demographics": {"age": "18-34"},
        },
        "white_collar": {
            "interests": ["Career development", "Tech", "Business", "Entrepreneurship"],
            "hashtags": ["#techjobs", "#careertok", "#jobhunt", "#corporatelife"],
            "demographics": {"age": "22-40"},
        },
    },
    "instagram": {
        "targeting_types": [
            "Interest",
            "Behavior",
            "Lookalike",
            "Custom Audience",
            "Location",
        ],
        "default": {
            "interests": ["Job searching", "Career", "Professional development"],
            "behaviors": ["Job seekers"],
            "demographics": {"age": "18-45"},
        },
    },
    "google_ads": {
        "targeting_types": [
            "Keywords",
            "In-Market Audiences",
            "Affinity Audiences",
            "Remarketing",
            "Demographics",
        ],
        "keyword_strategy": {
            "branded": "Company + jobs",
            "role_specific": "[role title] + jobs + [location]",
            "industry": "[industry] + careers + hiring",
            "competitor": "Competitor company + careers",
        },
    },
}


# ═══════════════════════════════════════════════════════════════════════════════
# CONTENT CALENDAR TEMPLATES
# ═══════════════════════════════════════════════════════════════════════════════

CONTENT_TYPES: Dict[str, List[Dict[str, Any]]] = {
    "facebook": [
        {
            "type": "Job Posting",
            "frequency_per_week": 3,
            "format": "Image + Link",
            "engagement": "medium",
        },
        {
            "type": "Employee Spotlight",
            "frequency_per_week": 1,
            "format": "Video/Photo",
            "engagement": "high",
        },
        {
            "type": "Company Culture",
            "frequency_per_week": 1,
            "format": "Carousel",
            "engagement": "high",
        },
        {
            "type": "Hiring Event",
            "frequency_per_week": 0.5,
            "format": "Event/Image",
            "engagement": "medium",
        },
    ],
    "instagram": [
        {
            "type": "Day in the Life",
            "frequency_per_week": 2,
            "format": "Reels",
            "engagement": "high",
        },
        {
            "type": "Job Opening",
            "frequency_per_week": 2,
            "format": "Stories + Feed",
            "engagement": "medium",
        },
        {
            "type": "Team Feature",
            "frequency_per_week": 1,
            "format": "Carousel",
            "engagement": "high",
        },
        {
            "type": "Behind the Scenes",
            "frequency_per_week": 1,
            "format": "Stories",
            "engagement": "high",
        },
    ],
    "linkedin": [
        {
            "type": "Job Posting",
            "frequency_per_week": 3,
            "format": "Sponsored Content",
            "engagement": "medium",
        },
        {
            "type": "Thought Leadership",
            "frequency_per_week": 1,
            "format": "Article/Post",
            "engagement": "high",
        },
        {
            "type": "Employee Story",
            "frequency_per_week": 1,
            "format": "Video/Image",
            "engagement": "high",
        },
        {
            "type": "Industry Insight",
            "frequency_per_week": 1,
            "format": "Document/Post",
            "engagement": "medium",
        },
    ],
    "tiktok": [
        {
            "type": "Day in the Life",
            "frequency_per_week": 3,
            "format": "Short Video",
            "engagement": "very_high",
        },
        {
            "type": "Job Announcement",
            "frequency_per_week": 2,
            "format": "Trending Audio",
            "engagement": "high",
        },
        {
            "type": "Behind the Scenes",
            "frequency_per_week": 1,
            "format": "Raw Video",
            "engagement": "high",
        },
        {
            "type": "Q&A / Tips",
            "frequency_per_week": 1,
            "format": "Duet/Stitch",
            "engagement": "high",
        },
    ],
    "twitter_x": [
        {
            "type": "Job Share",
            "frequency_per_week": 5,
            "format": "Text + Link",
            "engagement": "low",
        },
        {
            "type": "Company News",
            "frequency_per_week": 2,
            "format": "Thread",
            "engagement": "medium",
        },
        {
            "type": "Industry Commentary",
            "frequency_per_week": 2,
            "format": "Text",
            "engagement": "medium",
        },
    ],
    "youtube": [
        {
            "type": "Company Overview",
            "frequency_per_week": 0.25,
            "format": "Long-form Video",
            "engagement": "high",
        },
        {
            "type": "Employee Testimonial",
            "frequency_per_week": 0.5,
            "format": "Interview Video",
            "engagement": "high",
        },
        {
            "type": "Job Walk-through",
            "frequency_per_week": 0.5,
            "format": "Shorts/Video",
            "engagement": "medium",
        },
    ],
    "reddit": [
        {
            "type": "AMA / Hiring Thread",
            "frequency_per_week": 0.5,
            "format": "Text Post",
            "engagement": "high",
        },
        {
            "type": "Job Share",
            "frequency_per_week": 1,
            "format": "Promoted Post",
            "engagement": "medium",
        },
    ],
    "snapchat": [
        {
            "type": "Job Story",
            "frequency_per_week": 3,
            "format": "Snap Ad",
            "engagement": "medium",
        },
        {
            "type": "Day at Work",
            "frequency_per_week": 2,
            "format": "Story",
            "engagement": "high",
        },
    ],
    "pinterest": [
        {
            "type": "Career Infographic",
            "frequency_per_week": 2,
            "format": "Pin",
            "engagement": "medium",
        },
        {
            "type": "Culture Board",
            "frequency_per_week": 1,
            "format": "Board/Pin",
            "engagement": "medium",
        },
    ],
}


# ═══════════════════════════════════════════════════════════════════════════════
# CORE FUNCTIONS
# ═══════════════════════════════════════════════════════════════════════════════


def _get_collar_type(role: str, industry: str = "") -> Dict[str, Any]:
    """Get collar classification for a role."""
    if _HAS_COLLAR_INTEL:
        try:
            return _collar_intel.classify_collar(role, industry)
        except Exception as e:
            logger.warning("Collar classification failed: %s", e)
    # Fallback
    role_lower = role.lower()
    blue_kw = [
        "driver",
        "warehouse",
        "mechanic",
        "construction",
        "cook",
        "server",
        "cleaner",
        "laborer",
        "factory",
        "welder",
        "plumber",
        "electrician",
        "technician",
        "forklift",
        "cdl",
        "delivery",
        "picker",
        "packer",
    ]
    for kw in blue_kw:
        if kw in role_lower:
            return {
                "collar_type": "blue_collar",
                "confidence": 0.70,
                "method": "keyword_fallback",
            }
    return {
        "collar_type": "white_collar",
        "confidence": 0.50,
        "method": "default_fallback",
    }


def _get_benchmark(
    platform_key: str, industry: str, location: str = "", collar_type: str = "mixed"
) -> Dict[str, float]:
    """Get CPC/CPM/CTR/CPA benchmarks for a platform."""
    if _HAS_TREND_ENGINE and platform_key:
        try:
            benchmarks = {}
            for metric in ("cpc", "cpm", "ctr", "cpa"):
                result = _trend_engine.get_trend(
                    platform=platform_key,
                    industry=industry,
                    metric=metric,
                    years_back=1,
                )
                benchmarks[f"avg_{metric}"] = (
                    result.get("history", [{}])[-1].get("value") or 0
                    if result.get("history")
                    else 0
                )
            # Also get latest year data directly
            plat = platform_key.lower().replace(" ", "_")
            all_trends = getattr(_trend_engine, "_ALL_TRENDS", {})
            plat_data = all_trends.get(plat, {})
            ind_data = plat_data.get(industry, plat_data.get("general_entry_level", {}))
            yr_data = ind_data.get(2025, ind_data.get(2024, {}))
            if yr_data:
                return {
                    "avg_cpc": yr_data.get("avg_cpc", benchmarks.get("avg_cpc", 1.50)),
                    "avg_cpm": yr_data.get("avg_cpm", benchmarks.get("avg_cpm", 8.00)),
                    "avg_ctr": yr_data.get("avg_ctr", benchmarks.get("avg_ctr", 0.02)),
                    "avg_cpa": yr_data.get("avg_cpa", benchmarks.get("avg_cpa", 20.00)),
                    "avg_cvr": yr_data.get("avg_cvr", 0.08),
                }
            if benchmarks.get("avg_cpc"):
                return benchmarks
        except Exception as e:
            logger.warning("Benchmark lookup failed for %s: %s", platform_key, e)
    return {
        "avg_cpc": 1.50,
        "avg_cpm": 8.00,
        "avg_ctr": 0.025,
        "avg_cpa": 20.00,
        "avg_cvr": 0.08,
    }


def analyze_social_fit(
    role: str, industry: str, collar_type: str, goals: List[str] = None
) -> List[Dict[str, Any]]:
    """Rank social platforms by fit for this role/industry/collar combination.

    Returns list of dicts sorted by fit_score descending, each containing:
        platform_key, name, fit_score (0-100), reasons, color, ad_formats
    """
    if not goals:
        goals = ["job_applications"]

    collar = collar_type.lower().replace("-", "_").replace(" ", "_")
    if collar not in COLLAR_PLATFORM_SCORES:
        collar = "white_collar"

    collar_scores = COLLAR_PLATFORM_SCORES[collar]
    industry_boosts = INDUSTRY_PLATFORM_BOOSTS.get(industry, {})

    results = []
    for key, platform in SOCIAL_PLATFORMS.items():
        base_score = collar_scores.get(key, 40)
        score = float(base_score)

        # Apply goal weights
        goal_mult = 1.0
        for goal in goals:
            g = goal.lower().replace(" ", "_")
            weights = GOAL_PLATFORM_WEIGHTS.get(g, {})
            goal_mult *= weights.get(key, 1.0)
        score *= goal_mult

        # Apply industry boost
        ind_boost = industry_boosts.get(key, 1.0)
        score *= ind_boost

        # Cap at 100
        score = min(100, round(score))

        reasons = []
        if collar in platform.get("best_for_collar") or []:
            reasons.append(f"Strong fit for {collar.replace('_', ' ')} roles")
        if ind_boost > 1.0:
            reasons.append(
                f"Industry boost for {INDUSTRY_LABEL_MAP.get(industry, industry)}"
            )
        if goal_mult > 1.1:
            reasons.append(f"Aligned with campaign goals")
        if score < 40:
            reasons.append("Lower priority for this combination")

        results.append(
            {
                "platform_key": key,
                "name": platform["name"],
                "fit_score": score,
                "color": platform["color"],
                "icon": platform["icon"],
                "ad_formats": platform["ad_formats"],
                "strengths": platform["strengths"],
                "weaknesses": platform["weaknesses"],
                "best_times": platform["best_times"],
                "demographics": platform["demographics"],
                "reasons": reasons,
            }
        )

    results.sort(key=lambda x: x["fit_score"], reverse=True)
    return results


def analyze_search_fit(
    role: str, industry: str, location: str = "", collar_type: str = "white_collar"
) -> List[Dict[str, Any]]:
    """Rank search platforms by fit for this role/industry/location."""
    collar = collar_type.lower().replace("-", "_").replace(" ", "_")
    if collar not in COLLAR_PLATFORM_SCORES:
        collar = "white_collar"

    collar_scores = COLLAR_PLATFORM_SCORES[collar]
    industry_boosts = INDUSTRY_PLATFORM_BOOSTS.get(industry, {})

    results = []
    for key, platform in SEARCH_PLATFORMS.items():
        base_score = collar_scores.get(key, 50)
        score = float(base_score)

        ind_boost = industry_boosts.get(key, 1.0)
        score *= ind_boost
        score = min(100, round(score))

        reasons = []
        if collar in platform.get("best_for_collar") or []:
            reasons.append(f"Strong fit for {collar.replace('_', ' ')} hiring")
        if ind_boost > 1.0:
            reasons.append(f"Industry advantage")

        results.append(
            {
                "platform_key": key,
                "name": platform["name"],
                "fit_score": score,
                "color": platform["color"],
                "icon": platform["icon"],
                "type": platform["type"],
                "ad_formats": platform["ad_formats"],
                "strengths": platform["strengths"],
                "weaknesses": platform["weaknesses"],
                "reasons": reasons,
            }
        )

    results.sort(key=lambda x: x["fit_score"], reverse=True)
    return results


def allocate_social_budget(
    total_budget: float,
    social_platforms: List[Dict],
    search_platforms: List[Dict],
    goals: List[str] = None,
) -> Dict[str, Any]:
    """Allocate budget across social + search platforms based on fit scores.

    Uses fit scores as weights. Platforms with score < 30 are excluded.
    Returns dict with per-platform allocation + metadata.
    """
    if not goals:
        goals = ["job_applications"]

    # Filter to platforms with meaningful fit
    eligible_social = [p for p in social_platforms if p["fit_score"] >= 30]
    eligible_search = [p for p in search_platforms if p["fit_score"] >= 30]

    # Limit to top platforms by budget size
    if total_budget < 5000:
        max_social, max_search = 2, 1
    elif total_budget < 15000:
        max_social, max_search = 3, 2
    elif total_budget < 50000:
        max_social, max_search = 4, 2
    elif total_budget < 150000:
        max_social, max_search = 5, 3
    else:
        max_social, max_search = 6, 3

    eligible_social = eligible_social[:max_social]
    eligible_search = eligible_search[:max_search]

    all_platforms = eligible_social + eligible_search
    if not all_platforms:
        all_platforms = social_platforms[:2] + search_platforms[:1]

    # Determine social vs search split based on goals
    social_pct = 0.50
    is_awareness = any(g in ("brand_awareness", "employer_brand") for g in goals)
    is_applications = "job_applications" in goals
    if is_awareness and not is_applications:
        social_pct = 0.70
    elif is_applications and not is_awareness:
        social_pct = 0.40
    elif is_awareness and is_applications:
        social_pct = 0.55

    social_budget = total_budget * social_pct
    search_budget = total_budget * (1 - social_pct)

    # Allocate within social
    social_total_score = sum(p["fit_score"] for p in eligible_social) or 1
    social_allocations = {}
    for p in eligible_social:
        share = p["fit_score"] / social_total_score
        amt = round(social_budget * share, 2)
        social_allocations[p["platform_key"]] = {
            "budget": amt,
            "pct_of_total": round(amt / total_budget * 100, 1),
            "name": p["name"],
            "color": p["color"],
            "fit_score": p["fit_score"],
            "channel_type": "social",
        }

    # Allocate within search
    search_total_score = sum(p["fit_score"] for p in eligible_search) or 1
    search_allocations = {}
    for p in eligible_search:
        share = p["fit_score"] / search_total_score
        amt = round(search_budget * share, 2)
        search_allocations[p["platform_key"]] = {
            "budget": amt,
            "pct_of_total": round(amt / total_budget * 100, 1),
            "name": p["name"],
            "color": p["color"],
            "fit_score": p["fit_score"],
            "channel_type": "search",
        }

    all_allocations = {**social_allocations, **search_allocations}

    return {
        "total_budget": total_budget,
        "social_budget": round(social_budget, 2),
        "search_budget": round(search_budget, 2),
        "social_pct": round(social_pct * 100, 1),
        "search_pct": round((1 - social_pct) * 100, 1),
        "allocations": all_allocations,
        "platform_count": len(all_allocations),
        "excluded_platforms": [
            p["name"]
            for p in social_platforms + search_platforms
            if p["platform_key"] not in all_allocations
        ],
    }


def get_platform_benchmarks(
    allocations: Dict[str, Dict],
    industry: str,
    location: str = "",
    collar_type: str = "mixed",
) -> Dict[str, Dict]:
    """Get CPC/CPM/CPA benchmarks for each allocated platform."""
    results = {}

    def _fetch(key, alloc):
        platform_def = {**SOCIAL_PLATFORMS, **SEARCH_PLATFORMS}.get(key, {})
        te_key = platform_def.get("trend_engine_key")
        if te_key:
            bench = _get_benchmark(te_key, industry, location, collar_type)
        else:
            bench = platform_def.get(
                "benchmarks_2025",
                {"avg_cpc": 1.00, "avg_cpm": 6.00, "avg_ctr": 0.015, "avg_cpa": 15.00},
            )
        return key, {
            "platform": alloc.get("name", key),
            "avg_cpc": bench.get("avg_cpc", 1.00),
            "avg_cpm": bench.get("avg_cpm", 6.00),
            "avg_ctr": bench.get("avg_ctr", 0.015),
            "avg_cpa": bench.get("avg_cpa", 15.00),
            "avg_cvr": bench.get("avg_cvr", 0.08),
        }

    with ThreadPoolExecutor(max_workers=6) as pool:
        futures = {pool.submit(_fetch, k, v): k for k, v in allocations.items()}
        for future in as_completed(futures):
            try:
                key, data = future.result()
                results[key] = data
            except Exception as e:
                k = futures[future]
                logger.warning("Benchmark fetch error for %s: %s", k, e)
                results[k] = {
                    "platform": k,
                    "avg_cpc": 1.50,
                    "avg_cpm": 8.00,
                    "avg_ctr": 0.02,
                    "avg_cpa": 20.00,
                    "avg_cvr": 0.08,
                }

    return results


def generate_audience_targeting(
    role: str, industry: str, platform_key: str, collar_type: str = "white_collar"
) -> Dict[str, Any]:
    """Generate audience targeting suggestions for a specific platform."""
    collar = collar_type.lower().replace("-", "_").replace(" ", "_")
    template = TARGETING_TEMPLATES.get(platform_key, {})
    targeting_types = template.get(
        "targeting_types", ["Interest", "Location", "Demographics"]
    )

    # Get collar-specific or default targeting
    collar_data = template.get(
        collar, template.get("default", template.get("white_collar", {}))
    )

    role_lower = role.lower()
    industry_label = INDUSTRY_LABEL_MAP.get(
        industry, industry.replace("_", " ").title()
    )

    # Build job title suggestions
    title_variations = [role]
    words = role.split()
    if len(words) >= 2:
        title_variations.append(f"Senior {role}")
        title_variations.append(f"{role} - {industry_label}")

    result = {
        "platform": platform_key,
        "targeting_types": targeting_types,
        "job_titles": title_variations[:5],
        "industry_context": industry_label,
        "collar_type": collar,
    }

    if collar_data:
        result.update(
            {
                "interests": collar_data.get("interests") or [],
                "behaviors": collar_data.get("behaviors") or [],
                "demographics": collar_data.get("demographics", {}),
                "hashtags": collar_data.get("hashtags") or [],
            }
        )

    # Platform-specific additions
    if platform_key == "linkedin":
        result["seniority_levels"] = collar_data.get(
            "seniority", ["Entry", "Senior", "Manager"]
        )
        result["functions"] = collar_data.get("functions", ["Operations"])
        result["company_sizes"] = collar_data.get(
            "company_size", ["201-500", "501-1000"]
        )
    elif platform_key == "google_ads":
        kw_strat = template.get("keyword_strategy", {})
        result["keyword_strategy"] = {
            "branded": kw_strat.get("branded", f"Company + {role} jobs"),
            "role_specific": f"{role} jobs near me, {role} hiring, {role} careers",
            "industry": f"{industry_label} jobs, {industry_label} careers hiring",
            "long_tail": f"best {role} jobs, {role} salary, how to become a {role}",
        }

    return result


def generate_content_calendar(
    platforms: List[str], budget: float, duration_weeks: int = 4
) -> List[Dict[str, Any]]:
    """Generate a content calendar for the campaign duration."""
    calendar = []
    days_of_week = [
        "Monday",
        "Tuesday",
        "Wednesday",
        "Thursday",
        "Friday",
        "Saturday",
        "Sunday",
    ]

    for week_num in range(1, duration_weeks + 1):
        week_plan = {
            "week": week_num,
            "theme": _get_week_theme(week_num, duration_weeks),
            "posts": [],
        }

        for platform in platforms:
            content_types = CONTENT_TYPES.get(platform, [])
            for ct in content_types:
                freq = ct["frequency_per_week"]
                if freq < 1:
                    # Post every N weeks
                    interval = round(1 / freq)
                    if week_num % interval != 1:
                        continue
                    freq = 1

                for i in range(int(freq)):
                    day_idx = (
                        hash(f"{platform}_{ct['type']}_{i}") % 5
                    )  # Weekday posting
                    if ct.get("engagement") == "very_high":
                        day_idx = hash(f"{platform}_{i}") % 7  # Include weekends

                    week_plan["posts"].append(
                        {
                            "platform": platform,
                            "platform_name": SOCIAL_PLATFORMS.get(
                                platform, SEARCH_PLATFORMS.get(platform, {})
                            ).get("name", platform),
                            "content_type": ct["type"],
                            "format": ct["format"],
                            "day": days_of_week[day_idx % 7],
                            "expected_engagement": ct.get("engagement", "medium"),
                        }
                    )

        # Sort posts by day order
        day_order = {d: i for i, d in enumerate(days_of_week)}
        week_plan["posts"].sort(key=lambda x: day_order.get(x["day"], 0))
        calendar.append(week_plan)

    return calendar


def _get_week_theme(week_num: int, total_weeks: int) -> str:
    """Get thematic focus for each campaign week."""
    if total_weeks <= 2:
        themes = ["Launch & Awareness", "Drive Applications"]
    elif total_weeks <= 4:
        themes = [
            "Launch & Brand Awareness",
            "Role Spotlight & Engagement",
            "Employee Stories & Culture",
            "Push Applications & Retarget",
        ]
    elif total_weeks <= 8:
        themes = [
            "Launch & Brand Awareness",
            "Role Spotlight",
            "Employee Stories",
            "Culture & Benefits",
            "Industry Expertise",
            "Community Engagement",
            "Application Push",
            "Final Push & Retarget",
        ]
    else:
        cycle = [
            "Brand Awareness",
            "Role Highlight",
            "Culture Showcase",
            "Employee Spotlight",
            "Application Drive",
        ]
        themes = [cycle[(i) % len(cycle)] for i in range(total_weeks)]

    idx = min(week_num - 1, len(themes) - 1)
    return themes[idx]


def generate_creative_briefs(
    platforms: List[str], role: str, industry: str, collar_type: str = "white_collar"
) -> List[Dict[str, Any]]:
    """Generate creative brief per platform with specs, copy guidance, and tone."""
    industry_label = INDUSTRY_LABEL_MAP.get(
        industry, industry.replace("_", " ").title()
    )
    collar = collar_type.lower().replace("_", " ").title()
    briefs = []

    for platform_key in platforms:
        platform = SOCIAL_PLATFORMS.get(
            platform_key, SEARCH_PLATFORMS.get(platform_key, {})
        )
        if not platform:
            continue

        name = platform.get("name", platform_key)
        specs = platform.get("image_specs", {})
        copy_limits = platform.get("copy_limits", {})

        # Determine tone
        if platform_key in ("linkedin",):
            tone = "Professional, authoritative, data-driven"
            cta = "Apply Now | Learn More | View Role"
        elif platform_key in ("tiktok", "snapchat"):
            tone = "Casual, authentic, energetic, Gen Z-friendly"
            cta = "Apply in Bio | Swipe Up | Link in Comments"
        elif platform_key in ("instagram",):
            tone = "Visual, aspirational, culture-focused"
            cta = "Link in Bio | Apply Now | DM Us"
        elif platform_key in ("reddit",):
            tone = "Genuine, transparent, community-minded (avoid corporate speak)"
            cta = "Check out the role | AMA about working here"
        elif platform_key in ("twitter_x",):
            tone = "Concise, timely, conversational"
            cta = "Apply here | We're hiring | Join our team"
        else:
            tone = "Engaging, clear, benefit-focused"
            cta = "Apply Now | Learn More | Join Us"

        # Sample copy
        headline_copy = f"We're Hiring: {role} - {industry_label}"
        if collar_type in ("blue_collar", "pink_collar"):
            body_copy = (
                f"Looking for your next opportunity? We're hiring {role}s "
                f"in {industry_label}. Competitive pay, great benefits, "
                f"and a team that has your back. Apply today!"
            )
        else:
            body_copy = (
                f"Join our {industry_label} team as a {role}. "
                f"Grow your career with industry leaders, innovative projects, "
                f"and a culture that values your expertise. See the role."
            )

        briefs.append(
            {
                "platform_key": platform_key,
                "platform_name": name,
                "color": platform.get("color", "#333"),
                "tone": tone,
                "cta_options": cta,
                "image_specs": specs,
                "copy_limits": copy_limits,
                "sample_headline": headline_copy,
                "sample_body": body_copy,
                "recommended_formats": platform.get("ad_formats") or [][:4],
                "creative_tips": _get_creative_tips(platform_key, collar_type),
            }
        )

    return briefs


def _get_creative_tips(platform_key: str, collar_type: str) -> List[str]:
    """Return platform-specific creative tips."""
    tips = {
        "facebook": [
            "Use real employee photos over stock images (2-3x higher engagement)",
            "Lead Gen Forms reduce friction -- pre-fill from FB profile",
            "Video ads get 20-30% lower CPA than static images",
            "Include salary/pay range in copy (increases applies 30%+)",
        ],
        "instagram": [
            "Reels get 2x the reach of feed posts in 2025-2026",
            "Use behind-the-scenes content for authenticity",
            "Carousel ads with employee stories outperform single images",
            "Stories with polls/questions drive higher engagement",
        ],
        "linkedin": [
            "Document ads (PDF carousel) get 3x more clicks than image ads",
            "Employee advocacy posts get 8x more engagement than company posts",
            "Include specific job title and level in the headline",
            "Thought leadership content builds pipeline over time",
        ],
        "tiktok": [
            "Authentic, unpolished content outperforms corporate production",
            "Use trending sounds -- increases reach by 2-5x",
            "Day-in-the-life format is the #1 recruiting content type",
            "First 3 seconds are critical -- hook immediately",
        ],
        "youtube": [
            "Employee testimonials (2-3 min) are the highest performing recruitment videos",
            "Shorts (< 60s) are essential for discovery in 2026",
            "Include clear CTA overlay with apply link",
            "Thumbnail quality directly impacts click-through rate",
        ],
        "twitter_x": [
            "Threads about company culture get higher engagement than single tweets",
            "Use relevant industry hashtags (limit to 2-3 per post)",
            "Respond to all replies quickly to boost visibility",
        ],
        "reddit": [
            "Authenticity is paramount -- redditors detect and reject corporate speak",
            "AMA format for hiring managers drives quality engagement",
            "Target relevant subreddits (r/jobs, industry-specific subs)",
        ],
        "google_ads": [
            "Use responsive search ads with 10+ headline variations",
            "Location-specific keywords lower CPC by 20-30%",
            "Negative keywords are critical -- exclude 'salary', 'review', 'glassdoor'",
            "Performance Max campaigns can reduce CPA by 15-25% vs manual",
        ],
        "indeed_sponsored": [
            "Sponsored jobs get 4.5x more visibility than free postings",
            "Optimize job titles for search (avoid internal jargon)",
            "Include salary range -- improves apply rate by 30%+",
        ],
    }
    return tips.get(
        platform_key,
        ["Use high-quality visuals", "Include clear CTA", "Test multiple variations"],
    )


def estimate_performance(
    allocations: Dict[str, Dict], benchmarks: Dict[str, Dict], duration_weeks: int = 4
) -> Dict[str, Any]:
    """Estimate reach, impressions, clicks, and applications per platform."""
    platform_estimates = {}
    totals = {
        "impressions": 0,
        "clicks": 0,
        "applications": 0,
        "reach": 0,
        "total_spend": 0,
        "avg_cpa": 0,
    }

    for key, alloc in allocations.items():
        budget = alloc.get("budget") or 0
        if budget <= 0:
            continue

        bench = benchmarks.get(key, {})
        cpc = bench.get("avg_cpc", 1.50)
        cpm = bench.get("avg_cpm", 8.00)
        ctr = bench.get("avg_ctr", 0.02)
        cpa = bench.get("avg_cpa", 20.00)
        cvr = bench.get("avg_cvr", 0.08)

        # Estimate metrics
        est_clicks = round(budget / cpc) if cpc > 0 else 0
        est_impressions = round(budget / cpm * 1000) if cpm > 0 else 0
        est_applications = round(budget / cpa) if cpa > 0 else 0
        est_reach = round(est_impressions * 0.65)  # ~65% unique reach
        effective_cpa = (
            round(budget / est_applications, 2) if est_applications > 0 else 0
        )

        platform_estimates[key] = {
            "platform": alloc.get("name", key),
            "budget": budget,
            "impressions": est_impressions,
            "reach": est_reach,
            "clicks": est_clicks,
            "applications": est_applications,
            "cpc": cpc,
            "cpm": cpm,
            "ctr": round(ctr * 100, 2),
            "cpa": effective_cpa,
            "cvr": round(cvr * 100, 2),
            "roi_score": (
                round(est_applications / (budget / 1000), 2) if budget > 0 else 0
            ),
        }

        totals["impressions"] += est_impressions
        totals["clicks"] += est_clicks
        totals["applications"] += est_applications
        totals["reach"] += est_reach
        totals["total_spend"] += budget

    totals["avg_cpa"] = (
        round(totals["total_spend"] / totals["applications"], 2)
        if totals["applications"] > 0
        else 0
    )
    totals["avg_cpc"] = (
        round(totals["total_spend"] / totals["clicks"], 2)
        if totals["clicks"] > 0
        else 0
    )
    totals["overall_ctr"] = (
        round(totals["clicks"] / totals["impressions"] * 100, 2)
        if totals["impressions"] > 0
        else 0
    )

    # Rank by ROI
    roi_ranking = sorted(
        platform_estimates.items(),
        key=lambda x: x[1].get("roi_score") or 0,
        reverse=True,
    )
    best_roi = roi_ranking[0][1]["platform"] if roi_ranking else "N/A"
    worst_roi = roi_ranking[-1][1]["platform"] if roi_ranking else "N/A"

    return {
        "platforms": platform_estimates,
        "totals": totals,
        "best_roi_platform": best_roi,
        "worst_roi_platform": worst_roi,
        "duration_weeks": duration_weeks,
    }


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN ORCHESTRATOR
# ═══════════════════════════════════════════════════════════════════════════════


def generate_social_media_plan(
    role: str,
    location: str = "",
    industry: str = "general_entry_level",
    budget: float = 50000,
    goals: List[str] = None,
    duration_weeks: int = 4,
) -> Dict[str, Any]:
    """Main entry point -- generates a complete social + search media plan.

    Args:
        role: Job title (e.g., "Registered Nurse", "Software Engineer")
        location: Target location (e.g., "New York, NY")
        industry: Industry key from shared_utils.INDUSTRY_LABEL_MAP
        budget: Total campaign budget in USD
        goals: List of campaign goals
        duration_weeks: Campaign duration in weeks (1-12)

    Returns:
        Complete plan dict with all sections.
    """
    if not goals:
        goals = ["job_applications"]

    # Normalize inputs
    budget = parse_budget(budget, default=50000.0)
    if location:
        location = standardize_location(location)
    duration_weeks = max(1, min(12, duration_weeks))
    industry_label = INDUSTRY_LABEL_MAP.get(
        industry, industry.replace("_", " ").title()
    )

    # Step 1: Collar classification
    collar_result = _get_collar_type(role, industry)
    collar_type = collar_result.get("collar_type", "white_collar")

    # Step 2: Platform analysis (concurrent)
    with ThreadPoolExecutor(max_workers=2) as pool:
        social_future = pool.submit(
            analyze_social_fit, role, industry, collar_type, goals
        )
        search_future = pool.submit(
            analyze_search_fit, role, industry, location, collar_type
        )
        social_ranked = social_future.result()
        search_ranked = search_future.result()

    # Step 3: Budget allocation
    budget_plan = allocate_social_budget(budget, social_ranked, search_ranked, goals)
    allocations = budget_plan["allocations"]

    # Step 4: Benchmarks
    benchmarks = get_platform_benchmarks(allocations, industry, location, collar_type)

    # Step 5: Performance estimates
    performance = estimate_performance(allocations, benchmarks, duration_weeks)

    # Step 6: Audience targeting (concurrent)
    allocated_keys = list(allocations.keys())
    targeting = {}
    with ThreadPoolExecutor(max_workers=4) as pool:
        futures = {
            pool.submit(generate_audience_targeting, role, industry, k, collar_type): k
            for k in allocated_keys
        }
        for f in as_completed(futures):
            key = futures[f]
            try:
                targeting[key] = f.result()
            except Exception as e:
                logger.warning("Targeting gen failed for %s: %s", key, e)
                targeting[key] = {"platform": key, "error": str(e)}

    # Step 7: Content calendar
    calendar = generate_content_calendar(allocated_keys, budget, duration_weeks)

    # Step 8: Creative briefs
    briefs = generate_creative_briefs(allocated_keys, role, industry, collar_type)

    # Assemble final plan
    plan = {
        "success": True,
        "generated_at": datetime.utcnow().isoformat() + "Z",
        "inputs": {
            "role": role,
            "location": location,
            "industry": industry,
            "industry_label": industry_label,
            "budget": budget,
            "goals": goals,
            "duration_weeks": duration_weeks,
        },
        "collar_analysis": collar_result,
        "social_platforms": social_ranked,
        "search_platforms": search_ranked,
        "budget_allocation": budget_plan,
        "benchmarks": benchmarks,
        "performance": performance,
        "audience_targeting": targeting,
        "content_calendar": calendar,
        "creative_briefs": briefs,
        "summary": {
            "total_budget": budget,
            "platform_count": budget_plan["platform_count"],
            "collar_type": collar_type,
            "est_total_reach": performance["totals"]["reach"],
            "est_total_impressions": performance["totals"]["impressions"],
            "est_total_clicks": performance["totals"]["clicks"],
            "est_total_applications": performance["totals"]["applications"],
            "est_avg_cpa": performance["totals"]["avg_cpa"],
            "best_roi_platform": performance["best_roi_platform"],
            "top_social": social_ranked[0]["name"] if social_ranked else "N/A",
            "top_search": search_ranked[0]["name"] if search_ranked else "N/A",
        },
    }

    return plan


# ═══════════════════════════════════════════════════════════════════════════════
# EXCEL EXPORT
# Sapphire Blue palette: #0F172A, #2563EB, #DBEAFE, Calibri, column B start
# ═══════════════════════════════════════════════════════════════════════════════

# Design tokens
_XL_NAVY = "0F172A"
_XL_SAPPHIRE = "2563EB"
_XL_BLUE_LIGHT = "DBEAFE"
_XL_BLUE_PALE = "EFF6FF"
_XL_WHITE = "FFFFFF"
_XL_GREEN = "16A34A"
_XL_GREEN_BG = "DCFCE7"
_XL_AMBER = "D97706"
_XL_AMBER_BG = "FEF3C7"
_XL_MUTED = "78716C"
_XL_WARM_GRAY = "E7E5E4"


def generate_social_plan_excel(plan: Dict[str, Any]) -> bytes:
    """Generate Excel export of the social media plan.

    4 sheets: Executive Summary, Platform Strategy, Content Calendar, Performance.
    Sapphire Blue palette, Calibri font, content starts at column B.
    """
    if not _HAS_OPENPYXL:
        raise ImportError("openpyxl is required for Excel export")

    wb = Workbook()

    # Style helpers
    def _font(size=10, bold=False, color=_XL_NAVY):
        return Font(name="Calibri", size=size, bold=bold, color=color)

    def _fill(color):
        return PatternFill(start_color=color, end_color=color, fill_type="solid")

    def _align(h="left", v="center", wrap=False):
        return Alignment(horizontal=h, vertical=v, wrap_text=wrap)

    thin_border = Border(
        left=Side(style="thin", color=_XL_WARM_GRAY),
        right=Side(style="thin", color=_XL_WARM_GRAY),
        top=Side(style="thin", color=_XL_WARM_GRAY),
        bottom=Side(style="thin", color=_XL_WARM_GRAY),
    )

    inputs = plan.get("inputs", {})
    summary = plan.get("summary", {})
    allocations = plan.get("budget_allocation", {}).get("allocations", {})
    benchmarks = plan.get("benchmarks", {})
    performance = plan.get("performance", {})
    briefs = plan.get("creative_briefs") or []
    calendar = plan.get("content_calendar") or []

    # ── Sheet 1: Executive Summary ──
    ws = wb.active
    ws.title = "Executive Summary"
    ws.sheet_properties.tabColor = _XL_NAVY
    ws.column_dimensions["A"].width = 3
    ws.column_dimensions["B"].width = 30
    ws.column_dimensions["C"].width = 22
    ws.column_dimensions["D"].width = 22
    ws.column_dimensions["E"].width = 22
    ws.column_dimensions["F"].width = 22

    row = 2
    ws.merge_cells(f"B{row}:F{row}")
    c = ws[f"B{row}"]
    c.value = "Social & Search Media Plan"
    c.font = _font(18, True, _XL_WHITE)
    c.fill = _fill(_XL_NAVY)
    c.alignment = _align("center")
    row += 1
    ws.merge_cells(f"B{row}:F{row}")
    c = ws[f"B{row}"]
    c.value = f"{inputs.get('role', 'N/A')} | {inputs.get('location', 'N/A')} | {inputs.get('industry_label', 'N/A')}"
    c.font = _font(11, False, _XL_WHITE)
    c.fill = _fill(_XL_SAPPHIRE)
    c.alignment = _align("center")

    row += 2
    # Key metrics
    metrics = [
        ("Total Budget", f"${summary.get('total_budget') or 0:,.0f}"),
        ("Platforms", str(summary.get("platform_count") or 0)),
        ("Collar Type", summary.get("collar_type") or "".replace("_", " ").title()),
        ("Duration", f"{inputs.get('duration_weeks', 4)} weeks"),
        (
            "Campaign Goals",
            ", ".join(g.replace("_", " ").title() for g in inputs.get("goals") or []),
        ),
    ]
    for label, value in metrics:
        ws[f"B{row}"].value = label
        ws[f"B{row}"].font = _font(10, True)
        ws[f"C{row}"].value = value
        ws[f"C{row}"].font = _font(10)
        row += 1

    row += 1
    # Performance projections
    ws.merge_cells(f"B{row}:F{row}")
    c = ws[f"B{row}"]
    c.value = "Projected Performance"
    c.font = _font(13, True, _XL_WHITE)
    c.fill = _fill(_XL_SAPPHIRE)
    row += 1

    perf_headers = ["Metric", "Projected Value"]
    for i, h in enumerate(perf_headers):
        col = get_column_letter(2 + i)
        ws[f"{col}{row}"].value = h
        ws[f"{col}{row}"].font = _font(10, True, _XL_WHITE)
        ws[f"{col}{row}"].fill = _fill(_XL_NAVY)
    row += 1

    totals = performance.get("totals", {})
    perf_rows = [
        ("Est. Total Reach", f"{totals.get('reach') or 0:,}"),
        ("Est. Total Impressions", f"{totals.get('impressions') or 0:,}"),
        ("Est. Total Clicks", f"{totals.get('clicks') or 0:,}"),
        ("Est. Total Applications", f"{totals.get('applications') or 0:,}"),
        ("Est. Avg CPA", f"${totals.get('avg_cpa') or 0:,.2f}"),
        ("Est. Avg CPC", f"${totals.get('avg_cpc') or 0:,.2f}"),
        ("Best ROI Platform", performance.get("best_roi_platform", "N/A")),
    ]
    for label, value in perf_rows:
        ws[f"B{row}"].value = label
        ws[f"B{row}"].font = _font(10)
        ws[f"B{row}"].border = thin_border
        ws[f"C{row}"].value = value
        ws[f"C{row}"].font = _font(10, True)
        ws[f"C{row}"].border = thin_border
        if "CPA" in label or "CPC" in label:
            ws[f"C{row}"].fill = _fill(_XL_BLUE_PALE)
        row += 1

    # ── Sheet 2: Platform Strategy ──
    ws2 = wb.create_sheet("Platform Strategy")
    ws2.sheet_properties.tabColor = _XL_SAPPHIRE
    ws2.column_dimensions["A"].width = 3
    for col_letter in ["B", "C", "D", "E", "F", "G", "H", "I"]:
        ws2.column_dimensions[col_letter].width = 18

    row = 2
    ws2.merge_cells(f"B{row}:I{row}")
    c = ws2[f"B{row}"]
    c.value = "Platform Budget Allocation & Benchmarks"
    c.font = _font(14, True, _XL_WHITE)
    c.fill = _fill(_XL_NAVY)
    c.alignment = _align("center")
    row += 2

    headers = [
        "Platform",
        "Type",
        "Budget",
        "% of Total",
        "Fit Score",
        "Avg CPC",
        "Avg CPM",
        "Avg CPA",
    ]
    for i, h in enumerate(headers):
        col = get_column_letter(2 + i)
        ws2[f"{col}{row}"].value = h
        ws2[f"{col}{row}"].font = _font(10, True, _XL_WHITE)
        ws2[f"{col}{row}"].fill = _fill(_XL_NAVY)
        ws2[f"{col}{row}"].alignment = _align("center")
    row += 1

    for key, alloc in allocations.items():
        bench = benchmarks.get(key, {})
        vals = [
            alloc.get("name", key),
            alloc.get("channel_type", "social").title(),
            f"${alloc.get('budget') or 0:,.0f}",
            f"{alloc.get('pct_of_total') or 0:.1f}%",
            str(alloc.get("fit_score") or 0),
            f"${bench.get('avg_cpc') or 0:.2f}",
            f"${bench.get('avg_cpm') or 0:.2f}",
            f"${bench.get('avg_cpa') or 0:.2f}",
        ]
        bg = _XL_WHITE if row % 2 == 0 else _XL_BLUE_PALE
        for i, v in enumerate(vals):
            col = get_column_letter(2 + i)
            ws2[f"{col}{row}"].value = v
            ws2[f"{col}{row}"].font = _font(10)
            ws2[f"{col}{row}"].fill = _fill(bg)
            ws2[f"{col}{row}"].border = thin_border
            ws2[f"{col}{row}"].alignment = _align("center")
        row += 1

    # ── Sheet 3: Content Calendar ──
    ws3 = wb.create_sheet("Content Calendar")
    ws3.sheet_properties.tabColor = "16A34A"
    ws3.column_dimensions["A"].width = 3
    for col_letter in ["B", "C", "D", "E", "F"]:
        ws3.column_dimensions[col_letter].width = 22

    row = 2
    ws3.merge_cells(f"B{row}:F{row}")
    c = ws3[f"B{row}"]
    c.value = "Content Calendar"
    c.font = _font(14, True, _XL_WHITE)
    c.fill = _fill(_XL_NAVY)
    c.alignment = _align("center")
    row += 2

    for week in calendar:
        ws3.merge_cells(f"B{row}:F{row}")
        c = ws3[f"B{row}"]
        c.value = f"Week {week['week']}: {week['theme']}"
        c.font = _font(11, True, _XL_WHITE)
        c.fill = _fill(_XL_SAPPHIRE)
        row += 1

        cal_headers = ["Platform", "Content Type", "Format", "Day", "Engagement"]
        for i, h in enumerate(cal_headers):
            col = get_column_letter(2 + i)
            ws3[f"{col}{row}"].value = h
            ws3[f"{col}{row}"].font = _font(9, True)
            ws3[f"{col}{row}"].fill = _fill(_XL_BLUE_LIGHT)
        row += 1

        for post in week.get("posts") or [][:20]:  # Limit per week
            vals = [
                post.get("platform_name") or "",
                post.get("content_type") or "",
                post.get("format") or "",
                post.get("day") or "",
                post.get("expected_engagement") or "",
            ]
            for i, v in enumerate(vals):
                col = get_column_letter(2 + i)
                ws3[f"{col}{row}"].value = v
                ws3[f"{col}{row}"].font = _font(9)
                ws3[f"{col}{row}"].border = thin_border
            row += 1
        row += 1

    # ── Sheet 4: Performance Projections ──
    ws4 = wb.create_sheet("Performance Projections")
    ws4.sheet_properties.tabColor = "D97706"
    ws4.column_dimensions["A"].width = 3
    for col_letter in ["B", "C", "D", "E", "F", "G", "H", "I", "J"]:
        ws4.column_dimensions[col_letter].width = 16

    row = 2
    ws4.merge_cells(f"B{row}:J{row}")
    c = ws4[f"B{row}"]
    c.value = "Performance Projections by Platform"
    c.font = _font(14, True, _XL_WHITE)
    c.fill = _fill(_XL_NAVY)
    c.alignment = _align("center")
    row += 2

    perf_headers = [
        "Platform",
        "Budget",
        "Impressions",
        "Reach",
        "Clicks",
        "Applications",
        "CPC",
        "CPA",
        "ROI Score",
    ]
    for i, h in enumerate(perf_headers):
        col = get_column_letter(2 + i)
        ws4[f"{col}{row}"].value = h
        ws4[f"{col}{row}"].font = _font(10, True, _XL_WHITE)
        ws4[f"{col}{row}"].fill = _fill(_XL_NAVY)
        ws4[f"{col}{row}"].alignment = _align("center")
    row += 1

    plat_perf = performance.get("platforms", {})
    for key, pp in plat_perf.items():
        vals = [
            pp.get("platform", key),
            f"${pp.get('budget') or 0:,.0f}",
            f"{pp.get('impressions') or 0:,}",
            f"{pp.get('reach') or 0:,}",
            f"{pp.get('clicks') or 0:,}",
            str(pp.get("applications") or 0),
            f"${pp.get('cpc') or 0:.2f}",
            f"${pp.get('cpa') or 0:.2f}",
            f"{pp.get('roi_score') or 0:.1f}",
        ]
        bg = _XL_WHITE if row % 2 == 0 else _XL_BLUE_PALE
        for i, v in enumerate(vals):
            col = get_column_letter(2 + i)
            ws4[f"{col}{row}"].value = v
            ws4[f"{col}{row}"].font = _font(10)
            ws4[f"{col}{row}"].fill = _fill(bg)
            ws4[f"{col}{row}"].border = thin_border
            ws4[f"{col}{row}"].alignment = _align("center")
        row += 1

    # Totals row
    row += 1
    total_vals = [
        "TOTAL",
        f"${totals.get('total_spend') or 0:,.0f}",
        f"{totals.get('impressions') or 0:,}",
        f"{totals.get('reach') or 0:,}",
        f"{totals.get('clicks') or 0:,}",
        str(totals.get("applications") or 0),
        f"${totals.get('avg_cpc') or 0:.2f}",
        f"${totals.get('avg_cpa') or 0:.2f}",
        "",
    ]
    for i, v in enumerate(total_vals):
        col = get_column_letter(2 + i)
        ws4[f"{col}{row}"].value = v
        ws4[f"{col}{row}"].font = _font(10, True, _XL_WHITE)
        ws4[f"{col}{row}"].fill = _fill(_XL_NAVY)

    # Footer on all sheets
    for ws_ref in [ws, ws2, ws3, ws4]:
        max_row = ws_ref.max_row + 2
        ws_ref.merge_cells(f"B{max_row}:F{max_row}")
        c = ws_ref[f"B{max_row}"]
        c.value = f"Generated by Joveo Media Plan Generator | {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}"
        c.font = _font(8, False, _XL_MUTED)

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ═══════════════════════════════════════════════════════════════════════════════
# PPT EXPORT
# Joveo branding: Port Gore #202058, Blue Violet #5A54BD, Downy #6BB3CD
# ═══════════════════════════════════════════════════════════════════════════════


def generate_social_plan_ppt(plan: Dict[str, Any]) -> bytes:
    """Generate PowerPoint export of the social media plan.

    5 slides: Title, Budget Allocation, Platform Benchmarks, Performance, Recommendations.
    Joveo brand palette.
    """
    if not _HAS_PPTX:
        raise ImportError("python-pptx is required for PPT export")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Colors
    NAVY = RGBColor(0x20, 0x20, 0x58)
    BLUE = RGBColor(0x5A, 0x54, 0xBD)
    TEAL = RGBColor(0x6B, 0xB3, 0xCD)
    PINK = RGBColor(0xB5, 0x66, 0x9C)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    OFF_WHITE = RGBColor(0xFF, 0xFD, 0xF9)
    MUTED = RGBColor(0x59, 0x67, 0x80)
    DARK_TEXT = RGBColor(0x20, 0x20, 0x58)

    inputs = plan.get("inputs", {})
    summary = plan.get("summary", {})
    allocations = plan.get("budget_allocation", {}).get("allocations", {})
    benchmarks_data = plan.get("benchmarks", {})
    performance = plan.get("performance", {})

    def _add_bg(slide, color=NAVY):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_text_box(
        slide,
        left,
        top,
        width,
        height,
        text,
        font_size=12,
        bold=False,
        color=WHITE,
        align=PP_ALIGN.LEFT,
    ):
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = "Inter"
        p.alignment = align
        return txBox

    # ── Slide 1: Title ──
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide1, NAVY)

    _add_text_box(slide1, 1, 0.8, 11, 0.6, "JOVEO", 14, True, TEAL)
    _add_text_box(
        slide1,
        1,
        1.8,
        11,
        1.2,
        "Social & Search Media Plan",
        36,
        True,
        WHITE,
        PP_ALIGN.LEFT,
    )
    _add_text_box(
        slide1,
        1,
        3.2,
        11,
        0.8,
        f"{inputs.get('role', 'N/A')} | {inputs.get('location', 'N/A')} | {inputs.get('industry_label', 'N/A')}",
        16,
        False,
        TEAL,
    )

    # Key stats row
    stats = [
        ("Budget", f"${summary.get('total_budget') or 0:,.0f}"),
        ("Platforms", str(summary.get("platform_count") or 0)),
        ("Est. Applications", str(summary.get("est_total_applications") or 0)),
        ("Est. Avg CPA", f"${summary.get('est_avg_cpa') or 0:,.2f}"),
        ("Duration", f"{inputs.get('duration_weeks', 4)} Weeks"),
    ]
    x_start = 1.0
    for i, (label, value) in enumerate(stats):
        x = x_start + i * 2.3
        shape = slide1.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(4.5),
            Inches(2.0),
            Inches(1.5),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = BLUE
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "Inter"
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(9)
        p2.font.color.rgb = RGBColor(0xDD, 0xDB, 0xFF)
        p2.font.name = "Inter"
        p2.alignment = PP_ALIGN.CENTER

    # ── Slide 2: Budget Allocation ──
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide2, OFF_WHITE)

    _add_text_box(slide2, 0.8, 0.5, 11, 0.6, "Budget Allocation", 28, True, NAVY)

    bp = plan.get("budget_allocation", {})
    _add_text_box(
        slide2,
        0.8,
        1.2,
        11,
        0.4,
        f"Social: {bp.get('social_pct', 50)}% (${bp.get('social_budget') or 0:,.0f})  |  "
        f"Search: {bp.get('search_pct', 50)}% (${bp.get('search_budget') or 0:,.0f})",
        12,
        False,
        MUTED,
    )

    y = 2.0
    x = 0.8
    col = 0
    platform_colors = [
        BLUE,
        TEAL,
        PINK,
        RGBColor(0xCE, 0x90, 0x47),
        RGBColor(0x33, 0x87, 0x21),
        RGBColor(0xDC, 0x26, 0x26),
    ]
    for i, (key, alloc) in enumerate(allocations.items()):
        cx = x + (col * 4.1)
        cy = y + (i // 3) * 1.8

        shape = slide2.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(cx),
            Inches(cy),
            Inches(3.8),
            Inches(1.5),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = NAVY
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = alloc.get("name", key)
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = platform_colors[i % len(platform_colors)]
        p.font.name = "Inter"

        p2 = tf.add_paragraph()
        p2.text = (
            f"${alloc.get('budget') or 0:,.0f}  ({alloc.get('pct_of_total') or 0:.1f}%)"
        )
        p2.font.size = Pt(11)
        p2.font.color.rgb = WHITE
        p2.font.name = "Inter"

        p3 = tf.add_paragraph()
        p3.text = f"Fit Score: {alloc.get('fit_score') or 0} | {alloc.get('channel_type', 'social').title()}"
        p3.font.size = Pt(9)
        p3.font.color.rgb = RGBColor(0x8C, 0x96, 0xA8)
        p3.font.name = "Inter"

        col = (col + 1) % 3

    # ── Slide 3: Performance Projections ──
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide3, NAVY)

    _add_text_box(slide3, 0.8, 0.5, 11, 0.6, "Performance Projections", 28, True, WHITE)

    # Table-like layout
    headers = [
        "Platform",
        "Budget",
        "Impressions",
        "Clicks",
        "Applications",
        "CPC",
        "CPA",
    ]
    y_start = 1.5
    col_width = 1.65
    x_start = 0.8

    for i, h in enumerate(headers):
        _add_text_box(
            slide3,
            x_start + i * col_width,
            y_start,
            col_width,
            0.4,
            h,
            10,
            True,
            TEAL,
            PP_ALIGN.CENTER,
        )

    plat_perf = performance.get("platforms", {})
    for j, (key, pp) in enumerate(plat_perf.items()):
        y = y_start + 0.5 + j * 0.45
        vals = [
            pp.get("platform", key),
            f"${pp.get('budget') or 0:,.0f}",
            f"{pp.get('impressions') or 0:,}",
            f"{pp.get('clicks') or 0:,}",
            str(pp.get("applications") or 0),
            f"${pp.get('cpc') or 0:.2f}",
            f"${pp.get('cpa') or 0:.2f}",
        ]
        for i, v in enumerate(vals):
            _add_text_box(
                slide3,
                x_start + i * col_width,
                y,
                col_width,
                0.4,
                v,
                9,
                False,
                WHITE,
                PP_ALIGN.CENTER,
            )

    # Totals
    totals = performance.get("totals", {})
    y_total = y_start + 0.5 + len(plat_perf) * 0.45 + 0.2
    total_vals = [
        "TOTAL",
        f"${totals.get('total_spend') or 0:,.0f}",
        f"{totals.get('impressions') or 0:,}",
        f"{totals.get('clicks') or 0:,}",
        str(totals.get("applications") or 0),
        f"${totals.get('avg_cpc') or 0:.2f}",
        f"${totals.get('avg_cpa') or 0:.2f}",
    ]
    for i, v in enumerate(total_vals):
        _add_text_box(
            slide3,
            x_start + i * col_width,
            y_total,
            col_width,
            0.4,
            v,
            10,
            True,
            TEAL,
            PP_ALIGN.CENTER,
        )

    # ── Slide 4: Recommendations ──
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide4, OFF_WHITE)

    _add_text_box(
        slide4, 0.8, 0.5, 11, 0.6, "Recommendations & Next Steps", 28, True, NAVY
    )

    collar = summary.get("collar_type", "white_collar")
    recs = [
        f"Top Social Platform: {summary.get('top_social', 'N/A')} -- prioritize budget allocation here",
        f"Top Search Platform: {summary.get('top_search', 'N/A')} -- capture high-intent job seekers",
        f"Best ROI Platform: {performance.get('best_roi_platform', 'N/A')} -- scale spending if CPA holds",
        f"Collar Strategy: {collar.replace('_', ' ').title()} roles -- {'use Facebook/TikTok for reach, Indeed for volume' if 'blue' in collar else 'lead with LinkedIn + Google Ads for quality'}",
        f"A/B test ad creative in weeks 1-2, then consolidate budget to top performers",
        f"Set up conversion tracking on all platforms before launch",
        f"Monitor CPA weekly -- pause platforms exceeding 2x target CPA",
        f"Use retargeting audiences after week 2 to re-engage visitors who did not apply",
    ]
    y = 1.5
    for rec in recs:
        _add_text_box(slide4, 1.0, y, 10, 0.4, f"  {rec}", 11, False, DARK_TEXT)
        y += 0.5

    # Footer
    _add_text_box(
        slide4,
        0.8,
        6.8,
        11,
        0.3,
        f"Generated by Joveo Media Plan Generator | {datetime.utcnow().strftime('%Y-%m-%d')}",
        8,
        False,
        MUTED,
        PP_ALIGN.CENTER,
    )

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
