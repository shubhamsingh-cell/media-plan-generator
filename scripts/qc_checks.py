#!/usr/bin/env python3
"""Canonical QC gate for the 2026-07-03 bundle-quality fix program.

Regenerates the US sample (render_sample_outputs.py) and NZ sample
(render_sample_nz.py) and runs assertions mapped to
BUNDLE_QC_FINDINGS_2026-07-03.json. Every workstream (S1-S6, O1-O3) must run
this and report PASS/FAIL deltas before committing -- it is the single source
of truth for "did my fix work / did I break something else."

Exit code 0 = no FAILs (WARN allowed, printed for awareness).
Usage: python3 scripts/qc_checks.py
"""

import re
import sys
from pathlib import Path

PROJECT_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(PROJECT_ROOT))

import openpyxl  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.util import Emu  # noqa: E402

OUT_US = PROJECT_ROOT / "tmp_render"
OUT_NZ = PROJECT_ROOT / "tmp_render_nz"

results = []  # (level, check, detail)


def report(level, check, detail=""):
    results.append((level, check, detail))


def regenerate():
    import subprocess

    for script, outdir in (
        ("scripts/render_sample_outputs.py", OUT_US),
        ("scripts/render_sample_nz.py", OUT_NZ),
    ):
        p = subprocess.run(
            [sys.executable, script, str(outdir)],
            cwd=PROJECT_ROOT,
            capture_output=True,
            text=True,
        )
        if p.returncode != 0 or "FAIL" in p.stdout:
            report("FAIL", f"generation:{script}", p.stdout + p.stderr)
        else:
            report("PASS", f"generation:{script}")


# ---------------------------------------------------------------------------
# PPTX checks
# ---------------------------------------------------------------------------


def check_pptx(path: Path, label: str, expect_currency: str, forbid_currency: str):
    if not path.exists():
        report("FAIL", f"{label}:pptx-missing", str(path))
        return
    prs = Presentation(str(path))
    W, H = prs.slide_width, prs.slide_height
    all_text = []
    offslide = overlap = subfont = 0
    for si, slide in enumerate(prs.slides, 1):
        boxes = []
        for sh in slide.shapes:
            try:
                l, t, w, h = sh.left, sh.top, sh.width, sh.height
            except Exception:
                continue
            if l is None:
                continue
            if (
                l < -W * 0.01
                or t < -H * 0.01
                or (l + w) > W * 1.01
                or (t + h) > H * 1.01
            ):
                offslide += 1
            if sh.has_text_frame and sh.text_frame.text.strip():
                boxes.append((l, t, w, h))
                all_text.append(sh.text_frame.text)
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        if r.font.size and r.font.size.pt < 8 and r.text.strip():
                            subfont += 1
        for i in range(len(boxes)):
            for j in range(i + 1, len(boxes)):
                l1, t1, w1, h1 = boxes[i]
                l2, t2, w2, h2 = boxes[j]
                ix = max(0, min(l1 + w1, l2 + w2) - max(l1, l2))
                iy = max(0, min(t1 + h1, t2 + h2) - max(t1, t2))
                inter = ix * iy
                a = min(w1 * h1, w2 * h2)
                if a and inter / a > 0.30:
                    overlap += 1
    joined = "\n".join(all_text)

    report(
        "FAIL" if offslide else "PASS",
        f"{label}:no-offslide-shapes",
        f"{offslide} found",
    )
    report(
        "FAIL" if overlap else "PASS", f"{label}:no-text-overlap", f"{overlap} found"
    )
    report("FAIL" if subfont else "PASS", f"{label}:no-sub8pt-text", f"{subfont} found")

    # QA-leak checks (S1)
    for banned in (
        "Creative QC:",
        "Grade F",
        "heavily estimated",
        "Minimal data available",
    ):
        if banned in joined:
            report(
                "FAIL", f"{label}:no-internal-qa-leak", f"found {banned!r} in deck text"
            )
    if not any(
        b in joined
        for b in (
            "Creative QC:",
            "Grade F",
            "heavily estimated",
            "Minimal data available",
        )
    ):
        report("PASS", f"{label}:no-internal-qa-leak")

    # currency checks (S3)
    if expect_currency and expect_currency not in joined:
        report(
            "WARN",
            f"{label}:expected-currency-present",
            f"{expect_currency!r} not found anywhere",
        )
    if forbid_currency:
        bare = re.findall(r"(?<![A-Za-z$])\$\d", joined)
        if bare:
            report(
                "FAIL" if forbid_currency == "strict" else "WARN",
                f"{label}:no-bare-usd-on-local-plan",
                f"{len(bare)} bare-$ occurrences",
            )
        else:
            report("PASS", f"{label}:no-bare-usd-on-local-plan")

    # implausible "live" stats over zero data -- best-effort text heuristic
    if re.search(r"[\d,]{4,}\s+active jobs", joined):
        report(
            "WARN",
            f"{label}:live-postings-stat-present",
            "verify against workbook Postings field",
        )


# ---------------------------------------------------------------------------
# XLSX checks
# ---------------------------------------------------------------------------


def check_xlsx(
    path: Path, label: str, expect_currency_word: str, forbid_usd_header: bool
):
    if not path.exists():
        report("FAIL", f"{label}:xlsx-missing", str(path))
        return
    wb = openpyxl.load_workbook(path, data_only=False)

    # Total-row formula column alignment (S2): every =SUM(<col><n1>:<col><n2>)
    # formula must reference ITS OWN column letter.
    bad_totals = []
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for c in row:
                if isinstance(c.value, str) and c.value.startswith("=SUM("):
                    m = re.match(r"=SUM\(([A-Z]+)\d+:([A-Z]+)\d+\)", c.value)
                    if m:
                        col_letter = c.coordinate.rstrip("0123456789")
                        if m.group(1) != col_letter or m.group(2) != col_letter:
                            bad_totals.append(f"{ws.title}!{c.coordinate}={c.value}")
    report(
        "FAIL" if bad_totals else "PASS",
        f"{label}:total-row-formula-alignment",
        "; ".join(bad_totals[:10]),
    )

    # text-truncation heuristic: cell string length exactly at a suspicious
    # fixed cap (100/120/etc) with no trailing punctuation/whitespace
    truncated = []
    misspell = []
    text_money_like = 0
    numeric_money_like = 0
    usd_header_hits = []
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for c in row:
                v = c.value
                if not isinstance(v, str):
                    continue
                if len(v) in (100, 110, 120, 150, 200) and v[-1] not in ".!?)\"' ":
                    truncated.append(f"{ws.title}!{c.coordinate} ({len(v)} chars)")
                if "pheonix" in v.lower():
                    misspell.append(f"{ws.title}!{c.coordinate}")
                if forbid_usd_header and re.search(r"\(USD\)|Budget \(\$\)", v):
                    usd_header_hits.append(f"{ws.title}!{c.coordinate}={v!r}")
                if re.match(r"^\$[\d,]+(\.\d+)?$", v.strip()):
                    text_money_like += 1
    report(
        "FAIL" if truncated else "PASS",
        f"{label}:no-midword-truncation",
        "; ".join(truncated[:10]),
    )
    report(
        "FAIL" if misspell else "PASS",
        f"{label}:no-pheonix-misspelling",
        "; ".join(misspell[:10]),
    )
    if forbid_usd_header:
        report(
            "FAIL" if usd_header_hits else "PASS",
            f"{label}:no-usd-header-on-local-plan",
            "; ".join(usd_header_hits[:10]),
        )
    report("INFO", f"{label}:text-typed-money-cells", str(text_money_like))


def main():
    print("=" * 70)
    print("Regenerating samples...")
    regenerate()

    print("Checking US sample (should be bare $, no NZD)...")
    check_pptx(
        OUT_US / "sample_plan.pptx", "us-pptx", expect_currency="$", forbid_currency=""
    )
    check_xlsx(
        OUT_US / "sample_plan.xlsx",
        "us-xlsx",
        expect_currency_word="",
        forbid_usd_header=False,
    )

    print("Checking NZ sample (should be NZ$, never bare $ on plan figures)...")
    check_pptx(
        OUT_NZ / "sample_plan_nz.pptx",
        "nz-pptx",
        expect_currency="NZ$",
        forbid_currency="warn",
    )
    check_xlsx(
        OUT_NZ / "sample_plan_nz.xlsx",
        "nz-xlsx",
        expect_currency_word="NZD",
        forbid_usd_header=True,
    )

    print("=" * 70)
    fails = [r for r in results if r[0] == "FAIL"]
    warns = [r for r in results if r[0] == "WARN"]
    for level, check, detail in results:
        marker = {"PASS": "  ok ", "FAIL": "FAIL ", "WARN": "warn ", "INFO": "info "}[
            level
        ]
        line = f"[{marker}] {check}"
        if detail and level in ("FAIL", "WARN"):
            line += f"  -- {detail[:200]}"
        print(line)
    print("=" * 70)
    print(
        f"{len(fails)} FAIL, {len(warns)} WARN, {len(results) - len(fails) - len(warns)} PASS/INFO"
    )
    return 1 if fails else 0


if __name__ == "__main__":
    sys.exit(main())
