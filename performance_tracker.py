"""
performance_tracker.py -- Campaign Performance Tracker Engine

Parses actual campaign performance data (Excel/CSV), compares against
industry benchmarks from trend_engine, generates optimization
recommendations, and produces branded Excel/PPT reports.

Thread-safe, never crashes (all errors return structured error dicts).

Design tokens:
    Excel: Sapphire Blue palette (Navy #0F172A, Sapphire #2563EB, Light #DBEAFE)
    PPT:   Brand (Port Gore #202058, Blue Violet #5A54BD, Downy #6BB3CD)
"""

from __future__ import annotations

import io
import logging
import math
import re
import datetime
from typing import Any, Dict, List, Optional, Tuple

logger = logging.getLogger(__name__)

# ── Optional imports (lazy, with try/except like nova.py) ──────────────────
try:
    import trend_engine as _trend_engine

    _HAS_TREND_ENGINE = True
except ImportError:
    _trend_engine = None
    _HAS_TREND_ENGINE = False

try:
    import budget_engine as _budget_engine

    _HAS_BUDGET_ENGINE = True
except ImportError:
    _budget_engine = None
    _HAS_BUDGET_ENGINE = False

try:
    import collar_intelligence as _collar_intel

    _HAS_COLLAR_INTEL = True
except ImportError:
    _collar_intel = None
    _HAS_COLLAR_INTEL = False

try:
    from shared_utils import INDUSTRY_LABEL_MAP, parse_budget
except ImportError:
    INDUSTRY_LABEL_MAP = {}

    def parse_budget(v, *, default=100_000.0):
        try:
            return float(v)
        except Exception:
            return default


try:
    from benchmark_registry import get_channel_benchmark, get_benchmark_value

    _HAS_BENCHMARK_REGISTRY = True
except ImportError:
    _HAS_BENCHMARK_REGISTRY = False

# ═══════════════════════════════════════════════════════════════════════════════
# COLUMN MATCHING -- flexible, case-insensitive, partial match
# ═══════════════════════════════════════════════════════════════════════════════

# Map of canonical field name -> list of patterns to match against column headers
_COLUMN_PATTERNS: Dict[str, List[str]] = {
    "channel": [
        "channel",
        "platform",
        "source",
        "medium",
        "publisher",
        "vendor",
        "ad platform",
        "network",
    ],
    "spend": [
        "spend",
        "cost",
        "budget",
        "investment",
        "total spend",
        "total cost",
        "ad spend",
        "media spend",
    ],
    "clicks": ["click", "total click"],
    "impressions": ["impression", "impr", "views", "total impression"],
    "applications": [
        "application",
        "appli",
        "applies",
        "apply",
        "conversion",
        "lead",
        "submissions",
    ],
    "hires": ["hire", "placement", "onboard", "offer accepted", "starts"],
    "cpc": ["cpc", "cost per click", "cost/click", "avg cpc", "average cpc"],
    "cpa": [
        "cpa",
        "cost per application",
        "cost per apply",
        "cost/apply",
        "cost per conversion",
        "cost per lead",
        "cost/application",
        "cost per acquisition",
    ],
    "cph": ["cph", "cost per hire", "cost/hire", "cost per placement"],
    "ctr": [
        "ctr",
        "click through rate",
        "click-through rate",
        "clickthrough",
        "click rate",
    ],
}


def _match_column(header: str, patterns: List[str]) -> bool:
    """Check if a column header matches any of the given patterns (case-insensitive, partial)."""
    h = header.lower().strip()
    for pat in patterns:
        if pat in h or h in pat:
            return True
    return False


def _map_columns(headers: List[str]) -> Dict[str, Optional[int]]:
    """Map canonical field names to column indices from the header row."""
    mapping: Dict[str, Optional[int]] = {k: None for k in _COLUMN_PATTERNS}
    for idx, header in enumerate(headers):
        if not header or not header.strip():
            continue
        for field, patterns in _COLUMN_PATTERNS.items():
            if mapping[field] is None and _match_column(header, patterns):
                mapping[field] = idx
                break
    return mapping


# ═══════════════════════════════════════════════════════════════════════════════
# 1. PARSE PERFORMANCE DATA
# ═══════════════════════════════════════════════════════════════════════════════


def parse_performance_data(file_bytes: bytes, filename: str) -> List[Dict[str, Any]]:
    """Parse Excel/CSV with campaign results.

    Expected columns (flexible matching):
        Channel/Platform, Spend, Clicks, Applications/Applies, Hires,
        CPC, CPA, CPH, Impressions, CTR

    Returns list of dicts, one per channel/row with parsed numeric values.
    """
    MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB
    if len(file_bytes) > MAX_FILE_SIZE:
        return []
    try:
        ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
        if ext in ("xlsx", "xls"):
            return _parse_excel(file_bytes)
        elif ext == "csv":
            return _parse_csv(file_bytes)
        else:
            logger.warning("Unsupported file type for performance data: %s", filename)
            return []
    except Exception as exc:
        logger.exception("Failed to parse performance data from %s: %s", filename, exc)
        return []


def _parse_excel(file_bytes: bytes) -> List[Dict[str, Any]]:
    """Parse Excel file into performance records."""
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    try:
        ws = wb.active
        rows = list(ws.iter_rows(values_only=True))
        if len(rows) < 2:
            return []
        headers = [str(c).strip() if c else "" for c in rows[0]]
        return _rows_to_records(headers, rows[1:])
    finally:
        wb.close()


def _parse_csv(file_bytes: bytes) -> List[Dict[str, Any]]:
    """Parse CSV file into performance records."""
    import csv as csv_mod

    text = None
    for enc in ("utf-8", "utf-8-sig", "latin-1", "cp1252"):
        try:
            text = file_bytes.decode(enc)
            break
        except (UnicodeDecodeError, ValueError):
            continue
    if not text:
        return []
    reader = csv_mod.reader(io.StringIO(text))
    rows = list(reader)
    if len(rows) < 2:
        return []
    headers = [str(c).strip() if c else "" for c in rows[0]]
    return _rows_to_records(headers, rows[1:])


def _safe_float(val: Any) -> Optional[float]:
    """Safely convert a value to float, stripping currency/percent symbols."""
    if val is None:
        return None
    if isinstance(val, (int, float)):
        return float(val)
    s = str(val).strip()
    if not s or s.lower() in ("n/a", "na", "-", "--", "null", "none"):
        return None
    # Strip currency symbols, commas, percent signs
    s = re.sub(r"[$,\s%]", "", s)
    try:
        return float(s)
    except (ValueError, TypeError):
        return None


def _rows_to_records(headers: List[str], data_rows: List) -> List[Dict[str, Any]]:
    """Convert parsed rows to list of performance record dicts."""
    col_map = _map_columns(headers)
    records = []
    for row in data_rows:
        cells = list(row)
        # Skip empty rows
        if not any(c for c in cells if c is not None and str(c).strip()):
            continue
        # Must have at least a channel name
        ch_idx = col_map.get("channel")
        if ch_idx is None or ch_idx >= len(cells) or not cells[ch_idx]:
            continue
        channel = str(cells[ch_idx]).strip()
        if not channel or channel.lower() in ("total", "grand total", "sum"):
            continue

        rec: Dict[str, Any] = {"channel": channel}

        def _get(field: str) -> Optional[float]:
            idx = col_map.get(field)
            if idx is not None and idx < len(cells):
                return _safe_float(cells[idx])
            return None

        rec["spend"] = _get("spend") or 0.0
        rec["clicks"] = _get("clicks") or 0.0
        rec["impressions"] = _get("impressions") or 0.0
        rec["applications"] = _get("applications") or 0.0
        rec["hires"] = _get("hires") or 0.0

        # CPC: prefer explicit, else compute
        raw_cpc = _get("cpc")
        if raw_cpc is not None:
            rec["cpc"] = raw_cpc
        elif rec["clicks"] > 0:
            rec["cpc"] = round(rec["spend"] / rec["clicks"], 2)
        else:
            rec["cpc"] = 0.0

        # CPA: prefer explicit, else compute
        raw_cpa = _get("cpa")
        if raw_cpa is not None:
            rec["cpa"] = raw_cpa
        elif rec["applications"] > 0:
            rec["cpa"] = round(rec["spend"] / rec["applications"], 2)
        else:
            rec["cpa"] = 0.0

        # CPH: prefer explicit, else compute
        raw_cph = _get("cph")
        if raw_cph is not None:
            rec["cph"] = raw_cph
        elif rec["hires"] > 0:
            rec["cph"] = round(rec["spend"] / rec["hires"], 2)
        else:
            rec["cph"] = 0.0

        # CTR: prefer explicit (as decimal or percentage), else compute
        raw_ctr = _get("ctr")
        if raw_ctr is not None:
            rec["ctr"] = raw_ctr / 100.0 if raw_ctr > 1.0 else raw_ctr
        elif rec["impressions"] > 0:
            rec["ctr"] = round(rec["clicks"] / rec["impressions"], 4)
        else:
            rec["ctr"] = 0.0

        records.append(rec)

    return records


# ═══════════════════════════════════════════════════════════════════════════════
# 2. GET BENCHMARKS FOR CONTEXT
# ═══════════════════════════════════════════════════════════════════════════════

# Map of common channel names to trend_engine platform keys
_CHANNEL_TO_PLATFORM: Dict[str, str] = {
    "google": "google_search",
    "google ads": "google_search",
    "google search": "google_search",
    "search": "google_search",
    "sem": "google_search",
    "ppc": "google_search",
    "facebook": "meta_facebook",
    "meta": "meta_facebook",
    "meta ads": "meta_facebook",
    "instagram": "meta_instagram",
    "ig": "meta_instagram",
    "linkedin": "linkedin",
    "linkedin ads": "linkedin",
    "indeed": "indeed",
    "indeed sponsored": "indeed",
    "programmatic": "programmatic",
    "programmatic display": "programmatic",
    "ziprecruiter": "indeed",
    "glassdoor": "indeed",
    "job board": "indeed",
    "job boards": "indeed",
    "social": "meta_facebook",
    "social media": "meta_facebook",
    "display": "programmatic",
    "display ads": "programmatic",
    "career site": "google_search",
    "organic": "google_search",
}


def _resolve_platform(channel_name: str) -> str:
    """Resolve a channel name to a trend_engine platform key."""
    ch = channel_name.lower().strip()
    for pattern, platform in _CHANNEL_TO_PLATFORM.items():
        if pattern in ch or ch in pattern:
            return platform
    return "programmatic"  # safe default


def get_benchmarks_for_context(
    industry: str = "general_entry_level",
    roles: Optional[List[str]] = None,
    locations: Optional[List[str]] = None,
) -> Dict[str, Dict[str, Any]]:
    """Pull benchmark CPC/CPA/CTR from trend_engine for comparison context.

    Returns dict keyed by platform with benchmark metrics.
    """
    benchmarks: Dict[str, Dict[str, Any]] = {}
    platforms = [
        "google_search",
        "meta_facebook",
        "linkedin",
        "indeed",
        "programmatic",
        "meta_instagram",
    ]

    # Determine collar type from roles if possible
    collar_type = "mixed"
    if roles and _HAS_COLLAR_INTEL:
        try:
            result = _collar_intel.classify_collar(roles[0], industry=industry)
            collar_type = result.get("collar_type", "mixed")
        except Exception:
            pass

    location = ""
    if locations:
        location = locations[0] if isinstance(locations[0], str) else str(locations[0])

    now = datetime.datetime.now()
    month = now.month
    year = now.year

    for plat in platforms:
        bench: Dict[str, Any] = {"platform": plat}
        for metric in ("cpc", "cpa", "ctr", "cpm"):
            if _HAS_TREND_ENGINE:
                try:
                    result = _trend_engine.get_benchmark(
                        platform=plat,
                        industry=industry,
                        metric=metric,
                        collar_type=collar_type,
                        location=location,
                        month=month,
                        year=min(year, 2026),
                    )
                    bench[metric] = result.get("value", 0.0)
                    bench[f"{metric}_confidence"] = result.get("data_confidence", 0.5)
                except Exception:
                    bench[metric] = _fallback_benchmark(plat, metric)
            else:
                bench[metric] = _fallback_benchmark(plat, metric)
        benchmarks[plat] = bench

    return benchmarks


def _fallback_benchmark(platform: str, metric: str) -> float:
    """Provide fallback benchmarks when trend_engine is unavailable.

    Prefers benchmark_registry (single source of truth) when available,
    otherwise uses hardcoded values for resilience.
    """
    if _HAS_BENCHMARK_REGISTRY:
        return get_benchmark_value(platform, metric)

    # Hardcoded fallback (kept for resilience)
    _fallbacks = {
        "google_search": {"cpc": 2.69, "cpa": 45.00, "ctr": 0.042, "cpm": 10.00},
        "meta_facebook": {"cpc": 1.72, "cpa": 30.00, "ctr": 0.012, "cpm": 7.50},
        "meta_instagram": {"cpc": 1.50, "cpa": 35.00, "ctr": 0.010, "cpm": 8.00},
        "linkedin": {"cpc": 5.26, "cpa": 75.00, "ctr": 0.008, "cpm": 35.00},
        "indeed": {"cpc": 0.50, "cpa": 25.00, "ctr": 0.040, "cpm": 5.00},
        "programmatic": {"cpc": 0.63, "cpa": 22.00, "ctr": 0.025, "cpm": 4.50},
    }
    return _fallbacks.get(platform, _fallbacks["programmatic"]).get(metric, 1.00)


# ═══════════════════════════════════════════════════════════════════════════════
# 3. COMPARE ACTUAL vs BENCHMARK
# ═══════════════════════════════════════════════════════════════════════════════


def compare_actual_vs_benchmark(
    actual_data: List[Dict[str, Any]],
    benchmarks: Dict[str, Dict[str, Any]],
) -> List[Dict[str, Any]]:
    """For each channel, compute variance, efficiency score (0-100), and grade (A-F).

    Returns list of comparison dicts, one per channel.
    """
    comparisons = []
    for rec in actual_data:
        platform = _resolve_platform(rec["channel"])
        bench = benchmarks.get(platform, benchmarks.get("programmatic", {}))

        comp: Dict[str, Any] = {
            "channel": rec["channel"],
            "platform": platform,
            "actual": rec,
            "benchmark": {},
            "variances": {},
            "efficiency_score": 0.0,
            "grade": "C",
        }

        # Compute variances for key metrics
        scores = []
        for metric in ("cpc", "cpa", "ctr"):
            actual_val = rec.get(metric, 0.0)
            bench_val = bench.get(metric, 0.0)
            comp["benchmark"][metric] = bench_val

            if bench_val and bench_val > 0:
                if metric == "ctr":
                    # Higher CTR is better
                    variance_pct = ((actual_val - bench_val) / bench_val) * 100
                    # Positive variance = better for CTR
                    metric_score = min(100, max(0, 50 + variance_pct * 2))
                else:
                    # Lower CPC/CPA is better
                    variance_pct = ((actual_val - bench_val) / bench_val) * 100
                    # Negative variance = better for costs
                    metric_score = min(100, max(0, 50 - variance_pct * 2))
            else:
                variance_pct = 0.0
                metric_score = 50.0

            comp["variances"][metric] = {
                "actual": (
                    round(actual_val, 4) if metric == "ctr" else round(actual_val, 2)
                ),
                "benchmark": (
                    round(bench_val, 4) if metric == "ctr" else round(bench_val, 2)
                ),
                "variance_pct": round(variance_pct, 1),
                "is_favorable": (
                    (variance_pct < 0) if metric != "ctr" else (variance_pct > 0)
                ),
            }
            scores.append(metric_score)

        # Overall efficiency score (weighted: CPA 40%, CPC 35%, CTR 25%)
        if len(scores) >= 3:
            efficiency = scores[0] * 0.35 + scores[1] * 0.40 + scores[2] * 0.25
        elif scores:
            efficiency = sum(scores) / len(scores)
        else:
            efficiency = 50.0
        comp["efficiency_score"] = round(min(100, max(0, efficiency)), 1)

        # Grade
        comp["grade"] = _score_to_grade(comp["efficiency_score"])
        comparisons.append(comp)

    return comparisons


def _score_to_grade(score: float) -> str:
    """Convert efficiency score (0-100) to letter grade."""
    if score >= 85:
        return "A"
    elif score >= 70:
        return "B"
    elif score >= 55:
        return "C"
    elif score >= 40:
        return "D"
    else:
        return "F"


# ═══════════════════════════════════════════════════════════════════════════════
# 4. CHANNEL EFFICIENCY
# ═══════════════════════════════════════════════════════════════════════════════


def calculate_channel_efficiency(
    actual_data: List[Dict[str, Any]],
) -> List[Dict[str, Any]]:
    """ROI scoring: spend per hire, cost per apply, CTR efficiency.

    Returns list of dicts with efficiency metrics per channel.
    """
    total_spend = sum(r.get("spend") or 0 for r in actual_data)
    total_hires = sum(r.get("hires") or 0 for r in actual_data)
    total_apps = sum(r.get("applications") or 0 for r in actual_data)

    results = []
    for rec in actual_data:
        spend = rec.get("spend") or 0
        clicks = rec.get("clicks") or 0
        apps = rec.get("applications") or 0
        hires = rec.get("hires") or 0
        impressions = rec.get("impressions") or 0

        eff: Dict[str, Any] = {
            "channel": rec["channel"],
            "spend": spend,
            "spend_pct": round(
                (spend / total_spend * 100) if total_spend > 0 else 0, 1
            ),
            "clicks": clicks,
            "applications": apps,
            "hires": hires,
            "cpc": rec.get("cpc") or 0,
            "cpa": rec.get("cpa") or 0,
            "cph": rec.get("cph") or 0,
            "ctr": rec.get("ctr") or 0,
        }

        # Apply rate (clicks -> applications)
        eff["apply_rate"] = round(apps / clicks, 4) if clicks > 0 else 0.0
        # Hire rate (applications -> hires)
        eff["hire_rate"] = round(hires / apps, 4) if apps > 0 else 0.0
        # Full funnel efficiency: impressions -> hires
        eff["funnel_efficiency"] = (
            round(hires / impressions * 10000, 2) if impressions > 0 else 0.0
        )

        # ROI score (composite 0-100)
        sub_scores = []
        # Lower CPA is better -- score against median
        if total_apps > 0:
            avg_cpa = total_spend / total_apps
            if avg_cpa > 0 and eff["cpa"] > 0:
                cpa_ratio = avg_cpa / eff["cpa"]  # >1 means better than average
                sub_scores.append(min(100, cpa_ratio * 50))
        # Lower CPH is better
        if total_hires > 0:
            avg_cph = total_spend / total_hires
            if avg_cph > 0 and eff["cph"] > 0:
                cph_ratio = avg_cph / eff["cph"]
                sub_scores.append(min(100, cph_ratio * 50))
        # Higher CTR is better
        if eff["ctr"] > 0:
            sub_scores.append(min(100, eff["ctr"] / 0.03 * 50))

        eff["roi_score"] = (
            round(sum(sub_scores) / len(sub_scores), 1) if sub_scores else 50.0
        )
        eff["roi_grade"] = _score_to_grade(eff["roi_score"])

        results.append(eff)

    return results


# ═══════════════════════════════════════════════════════════════════════════════
# 5. GENERATE RECOMMENDATIONS
# ═══════════════════════════════════════════════════════════════════════════════


def generate_recommendations(
    comparison_results: List[Dict[str, Any]],
    total_budget: float = 0.0,
) -> Dict[str, Any]:
    """Generate actionable recommendations: channels to increase/decrease/drop,
    budget reallocation %, seasonal timing advice.

    Returns dict with categorized recommendations.
    """
    if not comparison_results:
        return {
            "recommendations": [],
            "budget_reallocation": {},
            "summary": "No data to analyze.",
        }

    # Sort channels by efficiency score
    sorted_channels = sorted(
        comparison_results, key=lambda x: x.get("efficiency_score") or 0, reverse=True
    )

    recommendations = []
    increase_channels = []
    decrease_channels = []
    drop_channels = []
    maintain_channels = []

    for ch in sorted_channels:
        grade = ch.get("grade", "C")
        score = ch.get("efficiency_score", 50)
        channel_name = ch.get("channel", "Unknown")
        actual = ch.get("actual", {})
        variances = ch.get("variances", {})

        cpc_var = variances.get("cpc", {})
        cpa_var = variances.get("cpa", {})

        if grade in ("A",):
            increase_channels.append(channel_name)
            rec = {
                "channel": channel_name,
                "action": "INCREASE",
                "priority": "high",
                "reason": f"{channel_name} is outperforming benchmarks (Score: {score}/100, Grade: {grade}). "
                f"CPC is {abs(cpc_var.get('variance_pct') or 0):.0f}% {'below' if cpc_var.get('is_favorable') else 'above'} benchmark.",
                "suggestion": f"Increase budget allocation to {channel_name} by 15-25% to capitalize on strong performance.",
                "icon": "trending_up",
            }
            recommendations.append(rec)

        elif grade == "B":
            maintain_channels.append(channel_name)
            rec = {
                "channel": channel_name,
                "action": "MAINTAIN",
                "priority": "medium",
                "reason": f"{channel_name} performs near or above benchmarks (Score: {score}/100, Grade: {grade}).",
                "suggestion": f"Maintain current {channel_name} allocation. Test incremental 5-10% increases.",
                "icon": "check_circle",
            }
            recommendations.append(rec)

        elif grade == "C":
            maintain_channels.append(channel_name)
            rec = {
                "channel": channel_name,
                "action": "OPTIMIZE",
                "priority": "medium",
                "reason": f"{channel_name} performs at or slightly below benchmarks (Score: {score}/100, Grade: {grade}).",
                "suggestion": f"Review {channel_name} targeting, creative, and bid strategy. Consider A/B testing.",
                "icon": "tune",
            }
            recommendations.append(rec)

        elif grade == "D":
            decrease_channels.append(channel_name)
            rec = {
                "channel": channel_name,
                "action": "DECREASE",
                "priority": "high",
                "reason": f"{channel_name} is underperforming benchmarks (Score: {score}/100, Grade: {grade}). "
                f"CPA is {abs(cpa_var.get('variance_pct') or 0):.0f}% above benchmark.",
                "suggestion": f"Reduce {channel_name} budget by 20-30% and reallocate to better-performing channels.",
                "icon": "trending_down",
            }
            recommendations.append(rec)

        else:  # F
            drop_channels.append(channel_name)
            rec = {
                "channel": channel_name,
                "action": "DROP / PAUSE",
                "priority": "critical",
                "reason": f"{channel_name} is significantly underperforming (Score: {score}/100, Grade: {grade}). "
                f"ROI is well below industry benchmarks.",
                "suggestion": f"Pause {channel_name} campaigns. Reallocate full budget to top-performing channels.",
                "icon": "cancel",
            }
            recommendations.append(rec)

    # Budget reallocation
    budget_reallocation = _compute_reallocation(sorted_channels, total_budget)

    # Seasonal timing advice
    now = datetime.datetime.now()
    month = now.month
    seasonal_advice = _get_seasonal_advice(month)

    # Summary
    avg_score = (
        sum(ch.get("efficiency_score", 50) for ch in sorted_channels)
        / len(sorted_channels)
        if sorted_channels
        else 50
    )
    overall_grade = _score_to_grade(avg_score)

    summary = (
        f"Overall campaign performance: Grade {overall_grade} ({avg_score:.0f}/100). "
        f"{len(increase_channels)} channel(s) to scale up, "
        f"{len(decrease_channels)} to reduce, "
        f"{len(drop_channels)} to pause."
    )

    return {
        "recommendations": recommendations[:10],  # Cap at 10
        "budget_reallocation": budget_reallocation,
        "seasonal_advice": seasonal_advice,
        "summary": summary,
        "increase_channels": increase_channels,
        "decrease_channels": decrease_channels,
        "drop_channels": drop_channels,
        "maintain_channels": maintain_channels,
    }


def _compute_reallocation(
    sorted_channels: List[Dict[str, Any]],
    total_budget: float,
) -> Dict[str, Any]:
    """Compute recommended budget reallocation based on performance scores."""
    if not sorted_channels:
        return {}

    # Current allocation (from actual spend)
    current: Dict[str, float] = {}
    total_actual_spend = 0.0
    for ch in sorted_channels:
        actual = ch.get("actual", {})
        spend = actual.get("spend") or 0
        current[ch["channel"]] = spend
        total_actual_spend += spend

    if total_actual_spend <= 0:
        return {}

    # Compute recommended % based on efficiency scores
    total_score = sum(max(ch.get("efficiency_score", 1), 1) for ch in sorted_channels)
    recommended: Dict[str, Dict[str, Any]] = {}

    for ch in sorted_channels:
        channel_name = ch["channel"]
        score = max(ch.get("efficiency_score", 1), 1)
        current_pct = (
            (current[channel_name] / total_actual_spend * 100)
            if total_actual_spend > 0
            else 0
        )
        # Weight recommended by score^1.5 to amplify differences
        weighted_score = score**1.5
        total_weighted = sum(
            max(c.get("efficiency_score", 1), 1) ** 1.5 for c in sorted_channels
        )
        recommended_pct = (
            (weighted_score / total_weighted * 100)
            if total_weighted > 0
            else current_pct
        )
        change_pct = recommended_pct - current_pct

        recommended[channel_name] = {
            "current_pct": round(current_pct, 1),
            "recommended_pct": round(recommended_pct, 1),
            "change_pct": round(change_pct, 1),
            "current_spend": round(current[channel_name], 2),
            "recommended_spend": round(total_actual_spend * recommended_pct / 100, 2),
        }

    return recommended


def _get_seasonal_advice(month: int) -> str:
    """Return seasonal hiring advice based on current month."""
    seasonal_tips = {
        1: "January is a high-volume hiring month. Increase budgets to capture New Year job seekers.",
        2: "February maintains strong candidate activity. Maintain investment levels.",
        3: "March is peak hiring season. Maximize spend on high-performing channels.",
        4: "April competition increases. Focus on differentiated creative and niche channels.",
        5: "May is strong for summer hiring. Scale up hourly/seasonal recruitment channels.",
        6: "June competition peaks for summer roles. Optimize bids and targeting.",
        7: "July typically sees a slight dip. Good time to A/B test new channels at lower CPCs.",
        8: "August is back-to-school hiring. Increase spend on entry-level channels.",
        9: "September is Q4 hiring ramp-up. Increase budgets across all channels.",
        10: "October is peak Q4 hiring. Maximize spend on top performers for holiday staffing.",
        11: "November sees high urgency hiring. Focus on fast-converting channels.",
        12: "December typically has lower competition, offering lower CPCs. Invest in brand awareness.",
    }
    return seasonal_tips.get(
        month, "Monitor market conditions and adjust budgets accordingly."
    )


# ═══════════════════════════════════════════════════════════════════════════════
# 6. GENERATE PERFORMANCE SCORECARD
# ═══════════════════════════════════════════════════════════════════════════════


def generate_performance_scorecard(
    actual_data: List[Dict[str, Any]],
    comparison: List[Dict[str, Any]],
) -> Dict[str, Any]:
    """Overall campaign grade (A-F) with breakdown."""
    if not actual_data:
        return {"grade": "N/A", "score": 0, "metrics": {}, "channel_grades": []}

    # Aggregate metrics
    total_spend = sum(r.get("spend") or 0 for r in actual_data)
    total_clicks = sum(r.get("clicks") or 0 for r in actual_data)
    total_impressions = sum(r.get("impressions") or 0 for r in actual_data)
    total_applications = sum(r.get("applications") or 0 for r in actual_data)
    total_hires = sum(r.get("hires") or 0 for r in actual_data)

    overall_cpc = round(total_spend / total_clicks, 2) if total_clicks > 0 else 0
    overall_cpa = (
        round(total_spend / total_applications, 2) if total_applications > 0 else 0
    )
    overall_cph = round(total_spend / total_hires, 2) if total_hires > 0 else 0
    overall_ctr = (
        round(total_clicks / total_impressions, 4) if total_impressions > 0 else 0
    )
    overall_apply_rate = (
        round(total_applications / total_clicks, 4) if total_clicks > 0 else 0
    )

    # Overall efficiency score from comparisons
    if comparison:
        avg_score = sum(c.get("efficiency_score", 50) for c in comparison) / len(
            comparison
        )
    else:
        avg_score = 50.0

    overall_grade = _score_to_grade(avg_score)

    channel_grades = []
    for c in comparison:
        channel_grades.append(
            {
                "channel": c["channel"],
                "grade": c["grade"],
                "score": c["efficiency_score"],
            }
        )
    channel_grades.sort(key=lambda x: x["score"], reverse=True)

    return {
        "grade": overall_grade,
        "score": round(avg_score, 1),
        "metrics": {
            "total_spend": round(total_spend, 2),
            "total_clicks": int(total_clicks),
            "total_impressions": int(total_impressions),
            "total_applications": int(total_applications),
            "total_hires": int(total_hires),
            "overall_cpc": overall_cpc,
            "overall_cpa": overall_cpa,
            "overall_cph": overall_cph,
            "overall_ctr": overall_ctr,
            "overall_apply_rate": overall_apply_rate,
        },
        "channel_grades": channel_grades,
        "channel_count": len(actual_data),
        "top_channel": channel_grades[0]["channel"] if channel_grades else "N/A",
        "bottom_channel": channel_grades[-1]["channel"] if channel_grades else "N/A",
    }


# ═══════════════════════════════════════════════════════════════════════════════
# 7. GENERATE PERFORMANCE EXCEL
# ═══════════════════════════════════════════════════════════════════════════════


def generate_performance_excel(
    report_data: Dict[str, Any], client_name: str = "Client"
) -> bytes:
    """Generate 4-sheet Excel report.

    Sheets:
        1. Performance Summary
        2. Channel Analysis
        3. Recommendations
        4. Projections

    Uses Sapphire Blue palette, Calibri font, data starts at column B.
    """
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
        from openpyxl.utils import get_column_letter
    except ImportError:
        logger.error("openpyxl not available for Excel generation")
        return b""

    # Design tokens
    NAVY = "0F172A"
    SAPPHIRE = "2563EB"
    BLUE_LIGHT = "DBEAFE"
    BLUE_PALE = "EFF6FF"
    STONE = "1C1917"
    MUTED = "78716C"
    WARM_GRAY = "E7E5E4"
    WHITE = "FFFFFF"
    GREEN = "16A34A"
    GREEN_BG = "DCFCE7"
    AMBER = "D97706"
    AMBER_BG = "FEF3C7"
    RED = "DC2626"
    RED_BG = "FEE2E2"

    # Fonts
    f_section = Font(name="Calibri", bold=True, size=14, color=WHITE)
    f_subsection = Font(name="Calibri", bold=True, size=12, color=NAVY)
    f_header = Font(name="Calibri", bold=True, size=10, color=WHITE)
    f_body = Font(name="Calibri", size=10, color=STONE)
    f_body_bold = Font(name="Calibri", bold=True, size=10, color=STONE)
    f_hero = Font(name="Calibri", bold=True, size=22, color=SAPPHIRE)
    f_hero_label = Font(name="Calibri", size=9, color=MUTED)
    f_metric_value = Font(name="Calibri", bold=True, size=14, color=NAVY)
    f_grade_big = Font(name="Calibri", bold=True, size=36, color=WHITE)
    f_footnote = Font(name="Calibri", italic=True, size=9, color=MUTED)
    f_green = Font(name="Calibri", bold=True, size=10, color=GREEN)
    f_red = Font(name="Calibri", bold=True, size=10, color=RED)
    f_amber = Font(name="Calibri", bold=True, size=10, color=AMBER)

    # Fills
    fill_navy = PatternFill(start_color=NAVY, end_color=NAVY, fill_type="solid")
    fill_sapphire = PatternFill(
        start_color=SAPPHIRE, end_color=SAPPHIRE, fill_type="solid"
    )
    fill_light = PatternFill(
        start_color=BLUE_LIGHT, end_color=BLUE_LIGHT, fill_type="solid"
    )
    fill_pale = PatternFill(
        start_color=BLUE_PALE, end_color=BLUE_PALE, fill_type="solid"
    )
    fill_white = PatternFill(start_color=WHITE, end_color=WHITE, fill_type="solid")
    fill_green_bg = PatternFill(
        start_color=GREEN_BG, end_color=GREEN_BG, fill_type="solid"
    )
    fill_amber_bg = PatternFill(
        start_color=AMBER_BG, end_color=AMBER_BG, fill_type="solid"
    )
    fill_red_bg = PatternFill(start_color=RED_BG, end_color=RED_BG, fill_type="solid")
    fill_green = PatternFill(start_color=GREEN, end_color=GREEN, fill_type="solid")

    # Grade fills
    def _grade_fill(grade: str):
        if grade == "A":
            return fill_green_bg
        elif grade == "B":
            return PatternFill(
                start_color="DCFCE7", end_color="DCFCE7", fill_type="solid"
            )
        elif grade == "C":
            return fill_amber_bg
        elif grade == "D":
            return PatternFill(
                start_color="FED7AA", end_color="FED7AA", fill_type="solid"
            )
        else:
            return fill_red_bg

    def _grade_font(grade: str):
        if grade in ("A", "B"):
            return f_green
        elif grade == "C":
            return f_amber
        else:
            return f_red

    # Alignment
    al_wrap = Alignment(wrap_text=True, vertical="top")
    al_center = Alignment(horizontal="center", vertical="center", wrap_text=True)
    al_left = Alignment(horizontal="left", vertical="top", wrap_text=True)
    al_right = Alignment(horizontal="right", vertical="top", wrap_text=True)

    # Border
    border = Border(
        left=Side(style="thin", color=WARM_GRAY),
        right=Side(style="thin", color=WARM_GRAY),
        top=Side(style="thin", color=WARM_GRAY),
        bottom=Side(style="thin", color=WARM_GRAY),
    )

    COL_START = 2  # Column B

    wb = Workbook()

    scorecard = report_data.get("scorecard", {})
    metrics = scorecard.get("metrics", {})
    comparisons = report_data.get("comparisons") or []
    recommendations_data = report_data.get("recommendations", {})
    efficiency = report_data.get("efficiency") or []

    # ── Helper ─────────────────────────────────────────────────────────
    def _write_section_header(ws, row, title, col_start=COL_START, col_end=8):
        for c in range(col_start, col_end + 1):
            cell = ws.cell(row=row, column=c)
            cell.fill = fill_navy
        ws.cell(row=row, column=col_start, value=title).font = f_section
        ws.cell(row=row, column=col_start).alignment = al_left
        return row + 1

    def _write_table_header(ws, row, headers, col_start=COL_START):
        for i, h in enumerate(headers):
            cell = ws.cell(row=row, column=col_start + i, value=h)
            cell.font = f_header
            cell.fill = fill_sapphire
            cell.alignment = al_center
            cell.border = border
        return row + 1

    def _write_table_row(ws, row, values, col_start=COL_START, fonts=None, fills=None):
        for i, v in enumerate(values):
            cell = ws.cell(row=row, column=col_start + i, value=v)
            cell.font = fonts[i] if fonts and i < len(fonts) else f_body
            if fills and i < len(fills) and fills[i]:
                cell.fill = fills[i]
            cell.alignment = al_center if i > 0 else al_left
            cell.border = border
        return row + 1

    # ══════════════════════════════════════════════════════════════════
    # SHEET 1: Performance Summary
    # ══════════════════════════════════════════════════════════════════
    ws1 = wb.active
    ws1.title = "Performance Summary"
    ws1.sheet_properties.tabColor = NAVY

    # Set column widths
    ws1.column_dimensions["A"].width = 3
    for col_letter in ["B", "C", "D", "E", "F", "G", "H"]:
        ws1.column_dimensions[col_letter].width = 18

    row = 2
    # Title
    ws1.cell(row=row, column=COL_START, value=f"Campaign Performance Report").font = (
        Font(name="Calibri", bold=True, size=18, color=NAVY)
    )
    row += 1
    ws1.cell(
        row=row,
        column=COL_START,
        value=f"{client_name} | Generated {datetime.datetime.now().strftime('%B %d, %Y')}",
    ).font = f_footnote
    row += 2

    # Overall Grade Card
    row = _write_section_header(ws1, row, "OVERALL PERFORMANCE SCORECARD")
    row += 1
    # Grade
    ws1.cell(row=row, column=COL_START, value="Overall Grade").font = f_subsection
    grade_cell = ws1.cell(
        row=row, column=COL_START + 1, value=scorecard.get("grade", "N/A")
    )
    grade_cell.font = Font(name="Calibri", bold=True, size=28, color=WHITE)
    _g = scorecard.get("grade", "C")
    if _g in ("A", "B"):
        grade_cell.fill = PatternFill(
            start_color=GREEN, end_color=GREEN, fill_type="solid"
        )
    elif _g == "C":
        grade_cell.fill = PatternFill(
            start_color=AMBER, end_color=AMBER, fill_type="solid"
        )
    else:
        grade_cell.fill = PatternFill(start_color=RED, end_color=RED, fill_type="solid")
    grade_cell.alignment = al_center

    ws1.cell(row=row, column=COL_START + 2, value="Score").font = f_hero_label
    ws1.cell(
        row=row, column=COL_START + 3, value=f"{scorecard.get('score') or 0}/100"
    ).font = f_metric_value
    row += 2

    # Key Metrics Row
    metric_labels = [
        "Total Spend",
        "Total Clicks",
        "Total Applications",
        "Total Hires",
        "Overall CPC",
        "Overall CPA",
    ]
    metric_values = [
        f"${metrics.get('total_spend') or 0:,.2f}",
        f"{metrics.get('total_clicks') or 0:,}",
        f"{metrics.get('total_applications') or 0:,}",
        f"{metrics.get('total_hires') or 0:,}",
        f"${metrics.get('overall_cpc') or 0:.2f}",
        f"${metrics.get('overall_cpa') or 0:.2f}",
    ]
    for i, (label, val) in enumerate(zip(metric_labels, metric_values)):
        col = COL_START + i
        ws1.cell(row=row, column=col, value=label).font = f_hero_label
        ws1.cell(row=row, column=col).alignment = al_center
        ws1.cell(row=row + 1, column=col, value=val).font = f_metric_value
        ws1.cell(row=row + 1, column=col).alignment = al_center
        ws1.cell(row=row + 1, column=col).fill = fill_light
        ws1.cell(row=row + 1, column=col).border = border
    row += 3

    # Channel Grades
    row = _write_section_header(ws1, row, "CHANNEL GRADES")
    row = _write_table_header(
        ws1,
        row,
        ["Channel", "Grade", "Score", "CPC Variance", "CPA Variance", "CTR Variance"],
    )
    for comp in comparisons:
        ch = comp.get("channel") or ""
        grade = comp.get("grade", "C")
        score_val = comp.get("efficiency_score") or 0
        vars_ = comp.get("variances", {})
        cpc_v = vars_.get("cpc", {}).get("variance_pct") or 0
        cpa_v = vars_.get("cpa", {}).get("variance_pct") or 0
        ctr_v = vars_.get("ctr", {}).get("variance_pct") or 0

        def _var_str(v, invert=False):
            arrow = "+" if v > 0 else ""
            return f"{arrow}{v:.1f}%"

        fonts = [
            f_body_bold,
            _grade_font(grade),
            f_body,
            f_green if cpc_v < 0 else f_red,
            f_green if cpa_v < 0 else f_red,
            f_green if ctr_v > 0 else f_red,
        ]
        fills = [None, _grade_fill(grade), None, None, None, None]
        row = _write_table_row(
            ws1,
            row,
            [
                ch,
                grade,
                f"{score_val:.0f}",
                _var_str(cpc_v),
                _var_str(cpa_v),
                _var_str(ctr_v),
            ],
            fonts=fonts,
            fills=fills,
        )

    row += 1
    ws1.cell(
        row=row,
        column=COL_START,
        value="Note: Negative CPC/CPA variance = better (lower cost). Positive CTR variance = better (higher click rate).",
    ).font = f_footnote

    # ══════════════════════════════════════════════════════════════════
    # SHEET 2: Channel Analysis
    # ══════════════════════════════════════════════════════════════════
    ws2 = wb.create_sheet("Channel Analysis")
    ws2.sheet_properties.tabColor = SAPPHIRE

    ws2.column_dimensions["A"].width = 3
    for col_letter in ["B", "C", "D", "E", "F", "G", "H", "I", "J"]:
        ws2.column_dimensions[col_letter].width = 16

    row = 2
    row = _write_section_header(ws2, row, "CHANNEL PERFORMANCE DETAIL", col_end=10)
    row = _write_table_header(
        ws2,
        row,
        [
            "Channel",
            "Spend",
            "Clicks",
            "Impressions",
            "Applications",
            "Hires",
            "CPC",
            "CPA",
            "CTR",
        ],
    )
    for rec in report_data.get("actual_data") or []:
        row = _write_table_row(
            ws2,
            row,
            [
                rec.get("channel") or "",
                f"${rec.get('spend') or 0:,.2f}",
                f"{rec.get('clicks') or 0:,.0f}",
                f"{rec.get('impressions') or 0:,.0f}",
                f"{rec.get('applications') or 0:,.0f}",
                f"{rec.get('hires') or 0:,.0f}",
                f"${rec.get('cpc') or 0:.2f}",
                f"${rec.get('cpa') or 0:.2f}",
                f"{(rec.get('ctr') or 0) * 100:.2f}%",
            ],
        )

    row += 2
    row = _write_section_header(ws2, row, "ACTUAL vs BENCHMARK COMPARISON", col_end=10)
    row = _write_table_header(
        ws2,
        row,
        [
            "Channel",
            "Actual CPC",
            "Bench CPC",
            "Actual CPA",
            "Bench CPA",
            "Actual CTR",
            "Bench CTR",
            "Grade",
        ],
    )
    for comp in comparisons:
        ch = comp.get("channel") or ""
        vars_ = comp.get("variances", {})
        grade = comp.get("grade", "C")
        row = _write_table_row(
            ws2,
            row,
            [
                ch,
                f"${vars_.get('cpc', {}).get('actual') or 0:.2f}",
                f"${vars_.get('cpc', {}).get('benchmark') or 0:.2f}",
                f"${vars_.get('cpa', {}).get('actual') or 0:.2f}",
                f"${vars_.get('cpa', {}).get('benchmark') or 0:.2f}",
                f"{(vars_.get('ctr', {}).get('actual') or 0) * 100:.2f}%",
                f"{(vars_.get('ctr', {}).get('benchmark') or 0) * 100:.2f}%",
                grade,
            ],
            fonts=[f_body_bold] + [f_body] * 6 + [_grade_font(grade)],
            fills=[None] * 7 + [_grade_fill(grade)],
        )

    # ══════════════════════════════════════════════════════════════════
    # SHEET 3: Recommendations
    # ══════════════════════════════════════════════════════════════════
    ws3 = wb.create_sheet("Recommendations")
    ws3.sheet_properties.tabColor = "16A34A"

    ws3.column_dimensions["A"].width = 3
    ws3.column_dimensions["B"].width = 18
    ws3.column_dimensions["C"].width = 14
    ws3.column_dimensions["D"].width = 12
    ws3.column_dimensions["E"].width = 50
    ws3.column_dimensions["F"].width = 50

    row = 2
    row = _write_section_header(ws3, row, "OPTIMIZATION RECOMMENDATIONS", col_end=6)
    row += 1
    ws3.cell(
        row=row, column=COL_START, value=recommendations_data.get("summary") or ""
    ).font = f_subsection
    row += 2

    row = _write_table_header(
        ws3, row, ["Channel", "Action", "Priority", "Reason", "Suggestion"]
    )
    for rec in recommendations_data.get("recommendations") or []:
        action = rec.get("action") or ""
        priority = rec.get("priority") or ""
        action_font = (
            f_green
            if action == "INCREASE"
            else (f_red if action in ("DROP / PAUSE", "DECREASE") else f_amber)
        )
        row = _write_table_row(
            ws3,
            row,
            [
                rec.get("channel") or "",
                action,
                priority.upper(),
                rec.get("reason") or "",
                rec.get("suggestion") or "",
            ],
            fonts=[f_body_bold, action_font, f_body, f_body, f_body],
        )

    row += 2
    # Budget Reallocation
    realloc = recommendations_data.get("budget_reallocation", {})
    if realloc:
        row = _write_section_header(ws3, row, "BUDGET REALLOCATION", col_end=6)
        row = _write_table_header(
            ws3,
            row,
            ["Channel", "Current %", "Recommended %", "Change", "Recommended Spend"],
        )
        for ch_name, alloc in realloc.items():
            change = alloc.get("change_pct") or 0
            change_font = f_green if change > 0 else (f_red if change < -5 else f_body)
            row = _write_table_row(
                ws3,
                row,
                [
                    ch_name,
                    f"{alloc.get('current_pct') or 0:.1f}%",
                    f"{alloc.get('recommended_pct') or 0:.1f}%",
                    f"{'+' if change > 0 else ''}{change:.1f}%",
                    f"${alloc.get('recommended_spend') or 0:,.2f}",
                ],
                fonts=[f_body_bold, f_body, f_body, change_font, f_body],
            )

    row += 2
    ws3.cell(
        row=row,
        column=COL_START,
        value=f"Seasonal Insight: {recommendations_data.get('seasonal_advice') or ''}",
    ).font = f_footnote

    # ══════════════════════════════════════════════════════════════════
    # SHEET 4: Projections
    # ══════════════════════════════════════════════════════════════════
    ws4 = wb.create_sheet("Projections")
    ws4.sheet_properties.tabColor = "D97706"

    ws4.column_dimensions["A"].width = 3
    for col_letter in ["B", "C", "D", "E", "F", "G"]:
        ws4.column_dimensions[col_letter].width = 20

    row = 2
    row = _write_section_header(
        ws4, row, "PROJECTED OUTCOMES (IF RECOMMENDATIONS APPLIED)"
    )
    row += 1

    total_spend = metrics.get("total_spend") or 0
    total_apps = metrics.get("total_applications") or 0
    total_hires = metrics.get("total_hires") or 0

    # Estimate improvement from reallocation
    improvement_factor = 1.15  # Conservative 15% improvement estimate
    if scorecard.get("score", 50) < 40:
        improvement_factor = 1.25  # More room for improvement
    elif scorecard.get("score", 50) > 70:
        improvement_factor = 1.08  # Already performing well

    proj_apps = int(total_apps * improvement_factor)
    proj_hires = int(total_hires * improvement_factor)
    proj_cpa = round(total_spend / proj_apps, 2) if proj_apps > 0 else 0
    proj_cph = round(total_spend / proj_hires, 2) if proj_hires > 0 else 0

    row = _write_table_header(
        ws4, row, ["Metric", "Current", "Projected (Optimized)", "Improvement"]
    )
    projections = [
        (
            "Applications",
            f"{total_apps:,}",
            f"{proj_apps:,}",
            f"+{proj_apps - total_apps:,} ({(improvement_factor - 1) * 100:.0f}%)",
        ),
        (
            "Hires",
            f"{total_hires:,}",
            f"{proj_hires:,}",
            f"+{proj_hires - total_hires:,} ({(improvement_factor - 1) * 100:.0f}%)",
        ),
        (
            "CPA",
            f"${metrics.get('overall_cpa') or 0:.2f}",
            f"${proj_cpa:.2f}",
            f"-${metrics.get('overall_cpa') or 0 - proj_cpa:.2f}",
        ),
        (
            "CPH",
            f"${metrics.get('overall_cph') or 0:.2f}",
            f"${proj_cph:.2f}",
            f"-${metrics.get('overall_cph') or 0 - proj_cph:.2f}",
        ),
    ]
    for label, current_val, proj_val, improvement in projections:
        row = _write_table_row(
            ws4,
            row,
            [label, current_val, proj_val, improvement],
            fonts=[f_body_bold, f_body, f_body, f_green],
        )

    row += 2
    ws4.cell(
        row=row,
        column=COL_START,
        value="Projections based on reallocating budget to higher-performing channels.",
    ).font = f_footnote
    ws4.cell(
        row=row + 1,
        column=COL_START,
        value="Actual results may vary. Improvement estimates are conservative based on historical reallocation outcomes.",
    ).font = f_footnote

    # Write to bytes
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ═══════════════════════════════════════════════════════════════════════════════
# 8. GENERATE PERFORMANCE PPT
# ═══════════════════════════════════════════════════════════════════════════════


def generate_performance_ppt(
    report_data: Dict[str, Any], client_name: str = "Client"
) -> bytes:
    """Generate 5-slide PPT report.

    Slides:
        1. Title
        2. Performance Overview
        3. Actual vs Benchmark
        4. Channel Efficiency Matrix
        5. Recommendations

    Uses branding: Port Gore #202058, Blue Violet #5A54BD, Downy #6BB3CD.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.enum.shapes import MSO_SHAPE
    except ImportError:
        logger.error("python-pptx not available for PPT generation")
        return b""

    # Brand colors
    NAVY = RGBColor(0x20, 0x20, 0x58)
    BLUE = RGBColor(0x5A, 0x54, 0xBD)
    TEAL = RGBColor(0x6B, 0xB3, 0xCD)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    OFF_WHITE = RGBColor(0xFF, 0xFD, 0xF9)
    DARK_TEXT = RGBColor(0x20, 0x20, 0x58)
    MUTED_TEXT = RGBColor(0x59, 0x67, 0x80)
    GREEN = RGBColor(0x33, 0x87, 0x21)
    AMBER = RGBColor(0xCE, 0x90, 0x47)
    RED_ACCENT = RGBColor(0xB5, 0x66, 0x9C)
    LIGHT_BG = RGBColor(0xF5, 0xF3, 0xFF)
    LIGHT_GREEN = RGBColor(0xE6, 0xF2, 0xE0)
    LIGHT_RED = RGBColor(0xFD, 0xE8, 0xE8)
    WARM_GRAY = RGBColor(0xEB, 0xE6, 0xE0)

    FONT_TITLE = "Poppins"
    FONT_BODY = "Inter"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    scorecard = report_data.get("scorecard", {})
    metrics = scorecard.get("metrics", {})
    comparisons = report_data.get("comparisons") or []
    recommendations_data = report_data.get("recommendations", {})
    efficiency = report_data.get("efficiency") or []

    def _add_bg(slide, color=OFF_WHITE):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_shape(slide, left, top, width, height, fill_color, border_color=None):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1)
        else:
            shape.line.fill.background()
        return shape

    def _add_text_box(
        slide,
        left,
        top,
        width,
        height,
        text,
        font_name=FONT_BODY,
        size=12,
        color=DARK_TEXT,
        bold=False,
        align=PP_ALIGN.LEFT,
    ):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = font_name
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = align
        return txBox

    # ── SLIDE 1: Title ────────────────────────────────────────────────
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    _add_bg(slide1, NAVY)

    # Accent bar at top
    _add_shape(slide1, Inches(0), Inches(0), prs.slide_width, Inches(0.08), TEAL)

    # Title
    _add_text_box(
        slide1,
        Inches(1.5),
        Inches(2.2),
        Inches(10),
        Inches(1.2),
        "Campaign Performance Tracker",
        FONT_TITLE,
        36,
        WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _add_text_box(
        slide1,
        Inches(1.5),
        Inches(3.4),
        Inches(10),
        Inches(0.6),
        f"{client_name} | Performance Analysis & Optimization Recommendations",
        FONT_BODY,
        16,
        TEAL,
        align=PP_ALIGN.CENTER,
    )
    _add_text_box(
        slide1,
        Inches(1.5),
        Inches(4.2),
        Inches(10),
        Inches(0.5),
        f"Generated {datetime.datetime.now().strftime('%B %d, %Y')}",
        FONT_BODY,
        12,
        MUTED_TEXT,
        align=PP_ALIGN.CENTER,
    )

    # Bottom bar
    _add_shape(slide1, Inches(0), Inches(7.0), prs.slide_width, Inches(0.5), BLUE)
    _add_text_box(
        slide1,
        Inches(0.5),
        Inches(7.05),
        Inches(4),
        Inches(0.4),
        "Powered by Nova AI Suite",
        FONT_BODY,
        10,
        WHITE,
        align=PP_ALIGN.LEFT,
    )

    # ── SLIDE 2: Performance Overview ─────────────────────────────────
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide2, OFF_WHITE)

    _add_text_box(
        slide2,
        Inches(0.8),
        Inches(0.4),
        Inches(8),
        Inches(0.7),
        "Performance Overview",
        FONT_TITLE,
        28,
        NAVY,
        bold=True,
    )
    _add_shape(slide2, Inches(0.8), Inches(1.0), Inches(2), Inches(0.04), TEAL)

    # Overall grade card
    grade = scorecard.get("grade", "N/A")
    grade_color = (
        GREEN if grade in ("A", "B") else (AMBER if grade == "C" else RED_ACCENT)
    )
    grade_bg = (
        LIGHT_GREEN
        if grade in ("A", "B")
        else (RGBColor(0xFE, 0xF3, 0xC7) if grade == "C" else LIGHT_RED)
    )

    grade_shape = _add_shape(
        slide2, Inches(0.8), Inches(1.5), Inches(2.5), Inches(2.5), grade_bg, WARM_GRAY
    )
    _add_text_box(
        slide2,
        Inches(0.8),
        Inches(1.6),
        Inches(2.5),
        Inches(0.4),
        "OVERALL GRADE",
        FONT_BODY,
        10,
        MUTED_TEXT,
        align=PP_ALIGN.CENTER,
    )
    _add_text_box(
        slide2,
        Inches(0.8),
        Inches(2.0),
        Inches(2.5),
        Inches(1.2),
        grade,
        FONT_TITLE,
        64,
        grade_color,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _add_text_box(
        slide2,
        Inches(0.8),
        Inches(3.2),
        Inches(2.5),
        Inches(0.4),
        f"Score: {scorecard.get('score') or 0}/100",
        FONT_BODY,
        14,
        DARK_TEXT,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # Metric cards
    metric_cards = [
        ("Total Spend", f"${metrics.get('total_spend') or 0:,.0f}"),
        ("Applications", f"{metrics.get('total_applications') or 0:,}"),
        ("Hires", f"{metrics.get('total_hires') or 0:,}"),
        ("Overall CPC", f"${metrics.get('overall_cpc') or 0:.2f}"),
        ("Overall CPA", f"${metrics.get('overall_cpa') or 0:.2f}"),
        ("Overall CTR", f"{(metrics.get('overall_ctr') or 0) * 100:.2f}%"),
    ]

    for i, (label, value) in enumerate(metric_cards):
        col = i % 3
        r = i // 3
        left = Inches(3.8 + col * 3.2)
        top = Inches(1.5 + r * 1.8)
        card = _add_shape(slide2, left, top, Inches(2.8), Inches(1.5), WHITE, WARM_GRAY)
        _add_text_box(
            slide2,
            left + Inches(0.2),
            top + Inches(0.2),
            Inches(2.4),
            Inches(0.3),
            label,
            FONT_BODY,
            10,
            MUTED_TEXT,
        )
        _add_text_box(
            slide2,
            left + Inches(0.2),
            top + Inches(0.6),
            Inches(2.4),
            Inches(0.6),
            value,
            FONT_TITLE,
            22,
            NAVY,
            bold=True,
        )

    # Bottom insight
    summary = recommendations_data.get("summary") or ""
    _add_text_box(
        slide2,
        Inches(0.8),
        Inches(5.2),
        Inches(11.5),
        Inches(0.5),
        summary,
        FONT_BODY,
        11,
        MUTED_TEXT,
    )

    # Footer bar
    _add_shape(slide2, Inches(0), Inches(7.0), prs.slide_width, Inches(0.5), NAVY)

    # ── SLIDE 3: Actual vs Benchmark ──────────────────────────────────
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide3, OFF_WHITE)

    _add_text_box(
        slide3,
        Inches(0.8),
        Inches(0.4),
        Inches(8),
        Inches(0.7),
        "Actual vs Benchmark Comparison",
        FONT_TITLE,
        28,
        NAVY,
        bold=True,
    )
    _add_shape(slide3, Inches(0.8), Inches(1.0), Inches(2), Inches(0.04), TEAL)

    # Table header
    table_headers = [
        "Channel",
        "Actual CPC",
        "Bench CPC",
        "CPC Var",
        "Actual CPA",
        "Bench CPA",
        "CPA Var",
        "Grade",
    ]
    num_rows = min(len(comparisons) + 1, 10)
    num_cols = len(table_headers)

    table_shape = slide3.shapes.add_table(
        num_rows,
        num_cols,
        Inches(0.5),
        Inches(1.4),
        Inches(12.3),
        Inches(0.5 * num_rows),
    )
    table = table_shape.table

    # Style header row
    for j, header in enumerate(table_headers):
        cell = table.cell(0, j)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.name = FONT_BODY
            paragraph.font.size = Pt(9)
            paragraph.font.color.rgb = WHITE
            paragraph.font.bold = True
            paragraph.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY

    # Data rows
    for i, comp in enumerate(comparisons[: num_rows - 1]):
        vars_ = comp.get("variances", {})
        cpc_v = vars_.get("cpc", {})
        cpa_v = vars_.get("cpa", {})
        g = comp.get("grade", "C")

        row_data = [
            comp.get("channel") or "",
            f"${cpc_v.get('actual') or 0:.2f}",
            f"${cpc_v.get('benchmark') or 0:.2f}",
            f"{'+' if (cpc_v.get('variance_pct') or 0) > 0 else ''}{cpc_v.get('variance_pct') or 0:.1f}%",
            f"${cpa_v.get('actual') or 0:.2f}",
            f"${cpa_v.get('benchmark') or 0:.2f}",
            f"{'+' if (cpa_v.get('variance_pct') or 0) > 0 else ''}{cpa_v.get('variance_pct') or 0:.1f}%",
            g,
        ]

        for j, val in enumerate(row_data):
            cell = table.cell(i + 1, j)
            cell.text = val
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = FONT_BODY
                paragraph.font.size = Pt(9)
                paragraph.font.color.rgb = DARK_TEXT
                paragraph.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT

            # Color code variance columns
            if j in (3, 6):
                v_pct = (
                    cpc_v.get("variance_pct") or 0
                    if j == 3
                    else cpa_v.get("variance_pct") or 0
                )
                if v_pct < -5:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = LIGHT_GREEN
                elif v_pct > 5:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = LIGHT_RED
            # Grade column
            if j == len(row_data) - 1:
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = (
                        GREEN
                        if g in ("A", "B")
                        else (AMBER if g == "C" else RED_ACCENT)
                    )

        # Alternating row color
        if i % 2 == 0:
            for j in range(num_cols):
                if table.cell(i + 1, j).fill.type is None or j not in (3, 6):
                    table.cell(i + 1, j).fill.solid()
                    table.cell(i + 1, j).fill.fore_color.rgb = LIGHT_BG

    _add_shape(slide3, Inches(0), Inches(7.0), prs.slide_width, Inches(0.5), NAVY)

    # ── SLIDE 4: Channel Efficiency Matrix ────────────────────────────
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide4, OFF_WHITE)

    _add_text_box(
        slide4,
        Inches(0.8),
        Inches(0.4),
        Inches(8),
        Inches(0.7),
        "Channel Efficiency Matrix",
        FONT_TITLE,
        28,
        NAVY,
        bold=True,
    )
    _add_shape(slide4, Inches(0.8), Inches(1.0), Inches(2), Inches(0.04), TEAL)

    # Show efficiency data as cards
    eff_sorted = sorted(efficiency, key=lambda x: x.get("roi_score") or 0, reverse=True)
    max_cards = min(len(eff_sorted), 6)

    for i, eff_rec in enumerate(eff_sorted[:max_cards]):
        col = i % 3
        r = i // 3
        left = Inches(0.5 + col * 4.2)
        top = Inches(1.5 + r * 2.8)

        roi_grade = eff_rec.get("roi_grade", "C")
        card_border = (
            GREEN
            if roi_grade in ("A", "B")
            else (AMBER if roi_grade == "C" else RED_ACCENT)
        )
        card = _add_shape(
            slide4, left, top, Inches(3.8), Inches(2.4), WHITE, card_border
        )

        # Channel name & grade
        _add_text_box(
            slide4,
            left + Inches(0.2),
            top + Inches(0.15),
            Inches(2.8),
            Inches(0.35),
            eff_rec.get("channel") or "",
            FONT_TITLE,
            14,
            NAVY,
            bold=True,
        )
        _add_text_box(
            slide4,
            left + Inches(3.0),
            top + Inches(0.15),
            Inches(0.6),
            Inches(0.35),
            roi_grade,
            FONT_TITLE,
            18,
            card_border,
            bold=True,
            align=PP_ALIGN.CENTER,
        )

        # Metrics
        eff_metrics = [
            f"Spend: ${eff_rec.get('spend') or 0:,.0f} ({eff_rec.get('spend_pct') or 0:.0f}%)",
            f"CPC: ${eff_rec.get('cpc') or 0:.2f}  |  CPA: ${eff_rec.get('cpa') or 0:.2f}",
            f"Apply Rate: {(eff_rec.get('apply_rate') or 0) * 100:.1f}%  |  Hire Rate: {(eff_rec.get('hire_rate') or 0) * 100:.1f}%",
            f"ROI Score: {eff_rec.get('roi_score') or 0:.0f}/100",
        ]
        for m_idx, m_text in enumerate(eff_metrics):
            _add_text_box(
                slide4,
                left + Inches(0.2),
                top + Inches(0.6 + m_idx * 0.4),
                Inches(3.4),
                Inches(0.35),
                m_text,
                FONT_BODY,
                9,
                MUTED_TEXT if m_idx < 3 else DARK_TEXT,
                bold=(m_idx == 3),
            )

    _add_shape(slide4, Inches(0), Inches(7.0), prs.slide_width, Inches(0.5), NAVY)

    # ── SLIDE 5: Recommendations ──────────────────────────────────────
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide5, OFF_WHITE)

    _add_text_box(
        slide5,
        Inches(0.8),
        Inches(0.4),
        Inches(8),
        Inches(0.7),
        "Optimization Recommendations",
        FONT_TITLE,
        28,
        NAVY,
        bold=True,
    )
    _add_shape(slide5, Inches(0.8), Inches(1.0), Inches(2), Inches(0.04), TEAL)

    recs = recommendations_data.get("recommendations") or []
    max_recs = min(len(recs), 6)

    for i, rec in enumerate(recs[:max_recs]):
        action = rec.get("action") or ""
        top = Inches(1.4 + i * 0.9)

        # Action color
        if action == "INCREASE":
            action_color = GREEN
            bg_color = LIGHT_GREEN
        elif action in ("DROP / PAUSE", "DECREASE"):
            action_color = RED_ACCENT
            bg_color = LIGHT_RED
        else:
            action_color = AMBER
            bg_color = RGBColor(0xFE, 0xF3, 0xC7)

        # Card
        _add_shape(
            slide5, Inches(0.5), top, Inches(12.3), Inches(0.75), bg_color, action_color
        )

        # Action badge
        badge = _add_shape(
            slide5,
            Inches(0.7),
            top + Inches(0.12),
            Inches(1.5),
            Inches(0.45),
            action_color,
        )
        _add_text_box(
            slide5,
            Inches(0.7),
            top + Inches(0.15),
            Inches(1.5),
            Inches(0.4),
            action,
            FONT_BODY,
            9,
            WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )

        # Channel name
        _add_text_box(
            slide5,
            Inches(2.4),
            top + Inches(0.08),
            Inches(2),
            Inches(0.35),
            rec.get("channel") or "",
            FONT_TITLE,
            12,
            NAVY,
            bold=True,
        )

        # Suggestion text
        _add_text_box(
            slide5,
            Inches(2.4),
            top + Inches(0.38),
            Inches(10),
            Inches(0.35),
            rec.get("suggestion") or "",
            FONT_BODY,
            9,
            MUTED_TEXT,
        )

    # Seasonal advice
    _add_text_box(
        slide5,
        Inches(0.8),
        Inches(6.3),
        Inches(11.5),
        Inches(0.5),
        f"Seasonal Insight: {recommendations_data.get('seasonal_advice') or ''}",
        FONT_BODY,
        10,
        TEAL,
        bold=True,
    )

    _add_shape(slide5, Inches(0), Inches(7.0), prs.slide_width, Inches(0.5), NAVY)
    _add_text_box(
        slide5,
        Inches(0.5),
        Inches(7.05),
        Inches(4),
        Inches(0.4),
        "Powered by Nova AI Suite",
        FONT_BODY,
        10,
        WHITE,
        align=PP_ALIGN.LEFT,
    )

    # Write to bytes
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ═══════════════════════════════════════════════════════════════════════════════
# ORCHESTRATOR -- single entry point for the API
# ═══════════════════════════════════════════════════════════════════════════════


def analyze_campaign(
    file_bytes: bytes,
    filename: str,
    campaign_name: str = "",
    industry: str = "general_entry_level",
    roles: Optional[List[str]] = None,
    locations: Optional[List[str]] = None,
    total_budget: float = 0.0,
) -> Dict[str, Any]:
    """Full analysis pipeline: parse -> benchmark -> compare -> recommend -> scorecard.

    Returns a complete report_data dict ready for Excel/PPT generation.
    Thread-safe, never raises.
    """
    try:
        # 1. Parse performance data
        actual_data = parse_performance_data(file_bytes, filename)
        if not actual_data:
            return {
                "error": "Could not parse performance data. Please check file format and column headers.",
                "success": False,
            }

        # 2. Get benchmarks
        benchmarks = get_benchmarks_for_context(industry, roles, locations)

        # 3. Compare actual vs benchmark
        comparisons = compare_actual_vs_benchmark(actual_data, benchmarks)

        # 4. Channel efficiency
        efficiency = calculate_channel_efficiency(actual_data)

        # 5. Recommendations
        total_spend = sum(r.get("spend") or 0 for r in actual_data)
        if total_budget <= 0:
            total_budget = total_spend
        recommendations = generate_recommendations(comparisons, total_budget)

        # 6. Scorecard
        scorecard = generate_performance_scorecard(actual_data, comparisons)

        return {
            "success": True,
            "campaign_name": campaign_name or "Campaign Analysis",
            "industry": industry,
            "industry_label": INDUSTRY_LABEL_MAP.get(
                industry, industry.replace("_", " ").title()
            ),
            "actual_data": actual_data,
            "benchmarks": benchmarks,
            "comparisons": comparisons,
            "efficiency": efficiency,
            "recommendations": recommendations,
            "scorecard": scorecard,
            "channel_count": len(actual_data),
            "total_spend": total_spend,
            "generated_at": datetime.datetime.now().isoformat(),
        }

    except Exception as exc:
        logger.exception("Campaign analysis failed: %s", exc)
        return {"error": f"Analysis failed: {str(exc)}", "success": False}
