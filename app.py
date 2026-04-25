#!/usr/bin/env python3
"""AI Media Planner - Standalone HTTP server with real research data."""

import html as html_mod
import json
import hmac
import os
import io
import datetime
import sys
import re
import zipfile
import uuid
import time
import secrets

try:
    import fcntl
except ImportError:
    fcntl = None  # Windows compatibility: file locking handled below
import gzip as gzip_module
import logging
import threading
import urllib.request
import urllib.error
from http.server import HTTPServer, BaseHTTPRequestHandler
from socketserver import ThreadingMixIn
import urllib.parse
from urllib.parse import urlparse
from pathlib import Path
from typing import Any, Optional
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter
from openpyxl.chart import Reference

from shared_utils import parse_budget

import benchmark_registry
import hashlib

# Static asset minifier -- minify CSS/JS at import time for zero-latency serving
try:
    from static_minifier import minify_static_assets, get_minified

    minify_static_assets()
except ImportError:
    get_minified = lambda _: None  # noqa: E731

logger = logging.getLogger(__name__)

# Auth module availability flag -- set during deferred_startup().
# Defaults to False so that if auth.py fails to import, admin endpoints are blocked.
_auth_module_loaded: bool = False


def _safe_str(val: object, default: str = "") -> str:
    """Safely convert any value to string, returning default for None/empty."""
    if val is None:
        return default
    if isinstance(val, str):
        return val
    return str(val)


def _internal_error_msg() -> str:
    """Return a user-facing error message with a short reference ID for support."""
    ref = uuid.uuid4().hex[:8]
    return f"Something went wrong (ref: {ref}). Please try again."


# ═══════════════════════════════════════════════════════════════════════════════
# WEB DATA ENRICHMENT (S72: Firecrawl removed -- module deleted)
# ═══════════════════════════════════════════════════════════════════════════════
# The firecrawl_enrichment module was deleted in S72 (credits exhausted, 0/500).
# The flag is kept as a permanent False so the dozen `if _firecrawl_available`
# guards scattered through the file fall through cleanly to their already-
# defined fallback paths (Tavily / Apify / Jina / Supabase / API integrations).
_firecrawl_available = False

# ═══════════════════════════════════════════════════════════════════════════════
# SENTRY WEBHOOK INTEGRATION (self-healing bridge)
# ═══════════════════════════════════════════════════════════════════════════════
try:
    from sentry_integration import (
        handle_sentry_webhook as _handle_sentry_webhook,
        get_sentry_status as _get_sentry_status,
        SentryAPIClient as _SentryAPIClient,
    )

    _sentry_integration_available = True
except ImportError:
    _sentry_integration_available = False

# ═══════════════════════════════════════════════════════════════════════════════
# API INTEGRATIONS MODULE (8 external data APIs for market intelligence)
# ═══════════════════════════════════════════════════════════════════════════════
try:
    from api_integrations import (
        fred as _api_fred,
        adzuna as _api_adzuna,
        jooble as _api_jooble,
        onet as _api_onet,
        bea as _api_bea,
        census as _api_census,
        usajobs as _api_usajobs,
        bls as _api_bls,
        test_all_apis as _api_test_all,
    )

    _api_integrations_available = True
except ImportError:
    _api_fred = _api_adzuna = _api_jooble = _api_onet = None
    _api_bea = _api_census = _api_usajobs = _api_bls = None
    _api_test_all = None
    _api_integrations_available = False
    logger.warning("api_integrations module not available; API enrichment disabled")

# ═══════════════════════════════════════════════════════════════════════════════
# ELEVENLABS INTEGRATION (TTS, STT, Sound Effects, Voice Design)
# ═══════════════════════════════════════════════════════════════════════════════
try:
    from elevenlabs_integration import (
        text_to_speech as _elevenlabs_tts,
        speech_to_text as _elevenlabs_stt,
        generate_sound_effect as _elevenlabs_sfx,
        generate_ad_voiceover as _elevenlabs_voiceover,
        check_elevenlabs_health as _elevenlabs_health,
    )

    _elevenlabs_available = True
except ImportError:
    _elevenlabs_available = False
    logger.warning(
        "elevenlabs_integration module not available; ElevenLabs features disabled"
    )

# ═══════════════════════════════════════════════════════════════════════════════
# SUPABASE DATA LAYER (Supabase-first with local JSON fallback)
# ═══════════════════════════════════════════════════════════════════════════════
try:
    from supabase_data import (
        get_knowledge,
        get_channel_benchmarks,
        get_salary_data,
        get_compliance_rules,
        get_market_trends,
        get_supply_repository,
        get_vendor_profiles,
    )

    _supabase_data_available = True
except ImportError:
    _supabase_data_available = False
    logger.warning("supabase_data module not available; Supabase enrichment disabled")

# ═══════════════════════════════════════════════════════════════════════════════
# ATTRIBUTION DASHBOARD (CFO-ready spend-to-hire attribution)
# ═══════════════════════════════════════════════════════════════════════════════
try:
    from attribution_dashboard import (
        generate_attribution_report as _attribution_report,
        generate_dashboard_html as _attribution_html,
    )

    _attribution_available = True
except ImportError:
    _attribution_available = False
    logger.warning(
        "attribution_dashboard module not available; attribution features disabled"
    )


# ═══════════════════════════════════════════════════════════════════════════════
# JOB SCRAPER (python-jobspy -- multi-platform job data)
# ═══════════════════════════════════════════════════════════════════════════════
try:
    from job_scraper import (
        scrape_jobs as _jobspy_scrape_jobs,
        get_job_market_stats as _jobspy_market_stats,
        get_salary_benchmarks as _jobspy_salary_benchmarks,
    )

    _jobspy_available = True
except ImportError:
    _jobspy_scrape_jobs = _jobspy_market_stats = _jobspy_salary_benchmarks = None
    _jobspy_available = False
    logger.warning("job_scraper module not available; JobSpy enrichment disabled")

# ═══════════════════════════════════════════════════════════════════════════════
# TAVILY AI SEARCH (AI-optimized web search)
# ═══════════════════════════════════════════════════════════════════════════════
try:
    from tavily_search import (
        search as _tavily_search,
        search_recruitment_news as _tavily_recruitment_news,
        research_company as _tavily_research_company,
    )

    _tavily_available = True
except ImportError:
    _tavily_search = _tavily_recruitment_news = _tavily_research_company = None
    _tavily_available = False
    logger.warning("tavily_search module not available; Tavily enrichment disabled")

# ═══════════════════════════════════════════════════════════════════════════════
# VECTOR SEARCH (Voyage AI embeddings + in-memory semantic search)
# ═══════════════════════════════════════════════════════════════════════════════
try:
    from vector_search import (
        # S61 FIX: use search_bounded (3s wall-clock cap) instead of raw
        # search().  The typo-hang RCA (claudedocs/s56_typo_hang_rca.md)
        # proved _enrich_chat_context was the bug: its
        # ThreadPoolExecutor(... ) as pool: context manager calls
        # shutdown(wait=True) on exit, which waits for the inner
        # _enrich_vector to finish. Raw search() can block up to 60s on
        # Voyage-AI rate-limit sleep, so the `fut.result(timeout=8)`
        # did nothing. search_bounded abandons via daemon thread.
        search_bounded as _vector_search,
        index_knowledge_base as _vector_index_kb,
    )

    _vector_search_available = True
except ImportError:
    _vector_search = _vector_index_kb = None
    _vector_search_available = False
    logger.warning("vector_search module not available; semantic search disabled")


# ═══════════════════════════════════════════════════════════════════════════════
# SLACK ALERTS MCP INTEGRATION (structured alerting for errors/deploys/health)
# ═══════════════════════════════════════════════════════════════════════════════
try:
    from slack_alerts import (
        send_alert as _slack_send_alert,
        send_error_alert as _slack_error_alert,
        get_status as _slack_alerts_status,
    )

    _slack_alerts_available = True
except ImportError:
    _slack_alerts_available = False
    logger.warning("slack_alerts module not available; Slack MCP alerting disabled")

# ═══════════════════════════════════════════════════════════════════════════════
# CALENDAR SYNC MCP INTEGRATION (Google Calendar for hiring milestones)
# ═══════════════════════════════════════════════════════════════════════════════
try:
    from calendar_sync import (
        create_hiring_event as _cal_create_event,
        get_upcoming_events as _cal_get_events,
        get_status as _cal_status,
    )

    _calendar_sync_available = True
except ImportError:
    _calendar_sync_available = False
    logger.warning("calendar_sync module not available; Calendar MCP disabled")

# ChromaDB RAG disabled -- module kept but disconnected from startup
_chroma_rag_available = False
_chroma_search = None
_chroma_stats = None

# ═══════════════════════════════════════════════════════════════════════════════
# NOVA INTELLIGENT CACHE (Supabase-backed query caching for instant responses)
# ═══════════════════════════════════════════════════════════════════════════════
try:
    from nova_cache import get_cache_stats as _nova_cache_stats

    _nova_cache_available = True
except ImportError:
    _nova_cache_available = False
    logger.warning("nova_cache module not available; intelligent caching disabled")

# ═══════════════════════════════════════════════════════════════════════════════
# POSTHOG ANALYTICS (stdlib-only, fire-and-forget event tracking)
# ═══════════════════════════════════════════════════════════════════════════════
try:
    from posthog_integration import (
        track_event as _ph_track,
        track_page_view as _ph_page_view,
        hash_ip as _ph_hash_ip,
        shutdown as _ph_shutdown,
    )

    _posthog_available = True
except ImportError:
    _posthog_available = False
    logger.warning("posthog_integration module not available; analytics disabled")


# ═══════════════════════════════════════════════════════════════════════════════
# LLM ROUTER (lazy import for narrative intelligence)
# ═══════════════════════════════════════════════════════════════════════════════
_llm_router_module = None
_llm_router_checked = False
_llm_router_lock = threading.Lock()


def _lazy_llm_router():
    """Lazy-load LLM router for narrative generation.

    Returns the module or None if unavailable. Thread-safe, loaded once.
    """
    global _llm_router_module, _llm_router_checked
    if _llm_router_checked:
        return _llm_router_module
    with _llm_router_lock:
        if _llm_router_checked:
            return _llm_router_module
        try:
            import llm_router as _mod

            _llm_router_module = _mod
        except ImportError:
            logger.warning("llm_router not available; narrative intelligence disabled")
            _llm_router_module = None
        _llm_router_checked = True
    return _llm_router_module


def _generate_narrative(
    data: dict, product_context: str, task_type: str = "", max_tokens: int = 300
) -> str:
    """Generate AI narrative for product output via LLM router.

    Routes all plan synthesis calls through TASK_CAMPAIGN_PLAN (which maps to
    Claude Haiku first) to ensure C-suite quality narrative output.

    Falls back to a template-based summary extracted from the data dict when
    all LLM providers are unavailable, ensuring the caller always gets usable
    text rather than an empty string.

    Args:
        data: Product data dict to summarize.
        product_context: Brief description of what this product does.
        task_type: LLM router task type (defaults to TASK_CAMPAIGN_PLAN for quality).
        max_tokens: Maximum tokens for the response.

    Returns:
        Generated narrative string, or template fallback on failure.
    """
    router = _lazy_llm_router()
    if not router:
        return _narrative_template_fallback(data, product_context)
    # S48: Use TASK_PLAN_NARRATIVE for prose generation -- routes to Groq (fast,
    # good narrative quality) instead of funneling through Gemini.  Falls back to
    # Gemini/Cerebras if Groq is unavailable.
    # S50: 10s timeout for all plan-gen LLM calls to avoid blocking the pipeline.
    task = task_type or getattr(router, "TASK_PLAN_NARRATIVE", "plan_narrative")
    try:
        prompt = (
            f"Based on this data, write a strategic recommendation as if presenting "
            f"to a VP of Talent Acquisition. Be specific with numbers, cite data points, "
            f"and explain the WHY behind each recommendation.\n\n"
            f"Data:\n{json.dumps(data, indent=2, default=str)[:2000]}\n\n"
            f"Context: {product_context}\n\n"
            f"Write 2-3 sentences. Include specific metrics. No generic filler."
        )
        result = router.call_llm(
            messages=[{"role": "user", "content": prompt}],
            system_prompt=(
                "You are a senior recruitment marketing strategist presenting to "
                "C-suite executives. Write with authority, cite specific data points, "
                "and explain causal reasoning. Every sentence must contain a number or "
                "a specific insight. No fluff, no platitudes."
            ),
            task_type=task,
            max_tokens=max_tokens,
            timeout_budget=10.0,
        )
        text = result.get("text") or ""
        if text:
            return text
        # LLM returned empty -- use template fallback
        return _narrative_template_fallback(data, product_context)
    except Exception as e:
        logger.error("Narrative generation failed: %s", e, exc_info=True)
        return _narrative_template_fallback(data, product_context)


def _narrative_template_fallback(data: dict, product_context: str) -> str:
    """Build a template-based narrative from raw data when LLM is unavailable.

    Extracts key numeric values from the data dict and formats them into a
    concise summary sentence. Returns empty string only if no useful data
    can be extracted.
    """
    parts: list = []
    # Extract common numeric fields from data
    if "roi_percentage" in data:
        roi = data.get("roi_percentage")
        payback = data.get("payback_period_months")
        savings = data.get("annual_savings")
        industry = data.get("industry") or "this sector"
        if roi:
            parts.append(f"Projected ROI of {roi}%")
        if payback:
            parts.append(f"payback period of {payback} months")
        if savings:
            try:
                parts.append(f"estimated annual savings of ${float(savings):,.0f}")
            except (ValueError, TypeError):
                parts.append(f"estimated annual savings of {savings}")
        if parts:
            return (
                f"For {industry} recruitment: {', '.join(parts)}. "
                f"See data sheets for detailed analysis."
            )
    # Generic fallback from product_context
    if product_context:
        return f"Analysis based on provided data. {product_context}"
    return ""


# ═══════════════════════════════════════════════════════════════════════════════
# LLM PRODUCT INSIGHTS -- optional AI insights for "dumb" products
# ═══════════════════════════════════════════════════════════════════════════════
import hashlib
from concurrent.futures import (
    ThreadPoolExecutor,
    as_completed,
)

_insights_cache: dict[str, tuple[str, float]] = {}
_insights_cache_lock = threading.Lock()
_INSIGHTS_CACHE_TTL = 3600.0  # 1 hour

# ═══════════════════════════════════════════════════════════════════════════════
# COPILOT SUGGESTION CACHE (60s TTL for inline suggestions)
# ═══════════════════════════════════════════════════════════════════════════════
_copilot_cache: dict[str, tuple[list, float]] = {}
_copilot_cache_lock = threading.Lock()
_COPILOT_CACHE_TTL = 60.0  # 60 seconds

# ═══════════════════════════════════════════════════════════════════════════════
# RESPONSE CACHE LIMITS & PERIODIC CLEANUP
# ═══════════════════════════════════════════════════════════════════════════════
_CACHE_MAX_ENTRIES = 500  # Max entries per cache before forced eviction
_CACHE_CLEANUP_INTERVAL = 300.0  # 5 minutes between cleanup sweeps
_FRAGMENT_CACHE_TTL = 600.0  # 10 minutes for HTML fragments (picks up file changes)


def _cache_cleanup_loop() -> None:
    """Periodically purge stale entries from all in-memory response caches.

    Runs every _CACHE_CLEANUP_INTERVAL seconds in a daemon thread.
    Removes expired entries and caps each cache at _CACHE_MAX_ENTRIES.
    """
    while True:
        try:
            time.sleep(_CACHE_CLEANUP_INTERVAL)
            now = time.time()

            # ── Insights cache ──
            try:
                with _insights_cache_lock:
                    stale_keys = [
                        k
                        for k, (_, ts) in _insights_cache.items()
                        if now - ts > _INSIGHTS_CACHE_TTL
                    ]
                    for k in stale_keys:
                        del _insights_cache[k]
                    # Cap size: evict oldest entries if still over limit
                    if len(_insights_cache) > _CACHE_MAX_ENTRIES:
                        sorted_keys = sorted(
                            _insights_cache, key=lambda k: _insights_cache[k][1]
                        )
                        for k in sorted_keys[
                            : len(_insights_cache) - _CACHE_MAX_ENTRIES
                        ]:
                            del _insights_cache[k]
                    purged_insights = len(stale_keys)
            except Exception as e:
                logger.error(
                    f"Cache cleanup failed for insights cache: {e}", exc_info=True
                )
                purged_insights = 0

            # ── Copilot cache ──
            try:
                with _copilot_cache_lock:
                    stale_keys = [
                        k
                        for k, (_, ts) in _copilot_cache.items()
                        if now - ts > _COPILOT_CACHE_TTL
                    ]
                    for k in stale_keys:
                        del _copilot_cache[k]
                    if len(_copilot_cache) > _CACHE_MAX_ENTRIES:
                        sorted_keys = sorted(
                            _copilot_cache, key=lambda k: _copilot_cache[k][1]
                        )
                        for k in sorted_keys[
                            : len(_copilot_cache) - _CACHE_MAX_ENTRIES
                        ]:
                            del _copilot_cache[k]
                    purged_copilot = len(stale_keys)
            except Exception as e:
                logger.error(
                    f"Cache cleanup failed for copilot cache: {e}", exc_info=True
                )
                purged_copilot = 0

            # ── Fragment cache (TTL-based invalidation for template reloads) ──
            try:
                with _fragment_cache_lock:
                    # Fragment cache stores str (no timestamp tuple), so track
                    # insertion time via a parallel dict if needed. For now,
                    # just cap the size since fragments are bounded by template count.
                    if len(_fragment_cache) > _CACHE_MAX_ENTRIES:
                        _fragment_cache.clear()
                        purged_fragments = 1
                    else:
                        purged_fragments = 0
            except Exception as e:
                logger.error(
                    f"Cache cleanup failed for fragment cache: {e}", exc_info=True
                )
                purged_fragments = 0

            # ── Shared plans store (24h TTL, max 1000) ──
            purged_shared_plans = 0
            try:
                with _shared_plans_lock:
                    stale_keys = [
                        k
                        for k, v in _shared_plans.items()
                        if now - v.get("created_at", 0) > _SHARED_PLANS_TTL
                    ]
                    for k in stale_keys:
                        del _shared_plans[k]
                    purged_shared_plans = len(stale_keys)
                    if len(_shared_plans) > _SHARED_PLANS_MAX:
                        sorted_keys = sorted(
                            _shared_plans,
                            key=lambda k: _shared_plans[k].get("created_at", 0),
                        )
                        for k in sorted_keys[: len(_shared_plans) - _SHARED_PLANS_MAX]:
                            del _shared_plans[k]
                            purged_shared_plans += 1
            except Exception as e:
                logger.error(
                    f"Cache cleanup failed for shared plans: {e}", exc_info=True
                )

            # ── Plan feedback store (24h TTL, max 1000 keys) ──
            purged_feedback = 0
            try:
                with _plan_feedback_lock:
                    stale_keys = [
                        k
                        for k, ts in _plan_feedback_ts.items()
                        if now - ts > _PLAN_FEEDBACK_TTL
                    ]
                    for k in stale_keys:
                        _plan_feedback.pop(k, None)
                        _plan_feedback_ts.pop(k, None)
                    purged_feedback = len(stale_keys)
                    if len(_plan_feedback) > _PLAN_FEEDBACK_MAX:
                        sorted_keys = sorted(
                            _plan_feedback_ts,
                            key=lambda k: _plan_feedback_ts.get(k, 0),
                        )
                        for k in sorted_keys[
                            : len(_plan_feedback) - _PLAN_FEEDBACK_MAX
                        ]:
                            _plan_feedback.pop(k, None)
                            _plan_feedback_ts.pop(k, None)
                            purged_feedback += 1
            except Exception as e:
                logger.error(
                    f"Cache cleanup failed for plan feedback: {e}", exc_info=True
                )

            # ── Scorecards store (12h TTL, max 500) ──
            purged_scorecards = 0
            try:
                with _scorecards_lock:
                    stale_keys = [
                        k
                        for k, (_, ts) in _scorecards.items()
                        if now - ts > _SCORECARDS_TTL
                    ]
                    for k in stale_keys:
                        del _scorecards[k]
                    purged_scorecards = len(stale_keys)
                    if len(_scorecards) > _SCORECARDS_MAX:
                        sorted_keys = sorted(
                            _scorecards,
                            key=lambda k: _scorecards[k][1],
                        )
                        for k in sorted_keys[: len(_scorecards) - _SCORECARDS_MAX]:
                            del _scorecards[k]
                            purged_scorecards += 1
            except Exception as e:
                logger.error(f"Cache cleanup failed for scorecards: {e}", exc_info=True)

            # ── Plan results store (30min TTL, background sweep) ──
            purged_plan_results = 0
            try:
                with _plan_results_lock:
                    stale_keys = [
                        k
                        for k, v in _plan_results_store.items()
                        if now - v.get("created", 0) > _PLAN_RESULTS_TTL_SECONDS
                    ]
                    for k in stale_keys:
                        del _plan_results_store[k]
                    purged_plan_results = len(stale_keys)
            except Exception as e:
                logger.error(
                    f"Cache cleanup failed for plan results: {e}", exc_info=True
                )

            total_purged = (
                purged_insights
                + purged_copilot
                + purged_fragments
                + purged_shared_plans
                + purged_feedback
                + purged_scorecards
                + purged_plan_results
            )
            if total_purged > 0:
                logger.info(
                    f"Cache cleanup: purged {purged_insights} insights, "
                    f"{purged_copilot} copilot, {purged_fragments} fragment, "
                    f"{purged_shared_plans} shared plans, {purged_feedback} feedback, "
                    f"{purged_scorecards} scorecards, {purged_plan_results} plan results"
                )
        except Exception as e:
            logger.error(f"Cache cleanup loop error: {e}", exc_info=True)


_cache_cleanup_thread = threading.Thread(
    target=_cache_cleanup_loop, daemon=True, name="response-cache-cleanup"
)
_cache_cleanup_thread.start()

# ═══════════════════════════════════════════════════════════════════════════════
# SLOW ENDPOINT PROFILER
# ═══════════════════════════════════════════════════════════════════════════════
_SLOW_THRESHOLD_MS = 1000.0  # Log endpoints taking >1 second
_SLOW_ENDPOINTS_MAX = 10  # Track top 10 slowest

_slow_endpoints_lock = threading.Lock()
# List of (path, method, duration_ms, timestamp) tuples, sorted by duration desc
_slow_endpoints: list[dict[str, Any]] = []


def _record_slow_endpoint(path: str, method: str, duration_ms: float) -> None:
    """Record an endpoint that exceeded the slow threshold.

    Maintains a sorted list of the top 10 slowest endpoints seen.
    Thread-safe.

    Args:
        path: The endpoint path (e.g., /api/chat).
        method: HTTP method (GET, POST).
        duration_ms: Request duration in milliseconds.
    """
    if duration_ms < _SLOW_THRESHOLD_MS:
        return

    logger.warning(f"SLOW ENDPOINT: {method} {path} took {duration_ms:.0f}ms")

    entry = {
        "path": path,
        "method": method,
        "duration_ms": round(duration_ms, 1),
        "timestamp": datetime.datetime.now(datetime.timezone.utc).isoformat(),
    }

    with _slow_endpoints_lock:
        _slow_endpoints.append(entry)
        # Sort by duration descending, keep top N
        _slow_endpoints.sort(key=lambda x: x["duration_ms"], reverse=True)
        del _slow_endpoints[_SLOW_ENDPOINTS_MAX:]


def get_slow_endpoints() -> list[dict[str, Any]]:
    """Return the top 10 slowest endpoints recorded.

    Returns:
        List of dicts with path, method, duration_ms, timestamp.
    """
    with _slow_endpoints_lock:
        return list(_slow_endpoints)


# ═══════════════════════════════════════════════════════════════════════════════
# BACKGROUND TASK EXECUTOR
# ═══════════════════════════════════════════════════════════════════════════════
_BG_DEFAULT_TIMEOUT = 30.0  # seconds

_background_executor = ThreadPoolExecutor(max_workers=8, thread_name_prefix="bg-task")
_bg_task_metrics_lock = threading.Lock()
_bg_task_metrics: dict[str, int] = {
    "submitted": 0,
    "completed": 0,
    "failed": 0,
    "timed_out": 0,
}


def submit_background_task(
    fn: Any,
    *args: Any,
    task_name: str = "",
    timeout: float = _BG_DEFAULT_TIMEOUT,
    **kwargs: Any,
) -> None:
    """Submit a function to run in the background thread pool.

    Wraps the function with error handling, timeout enforcement,
    and start/end logging. Never blocks the caller.

    Args:
        fn: Callable to execute in the background.
        *args: Positional arguments for fn.
        task_name: Human-readable name for logging (defaults to fn.__name__).
        timeout: Max seconds the task may run (default 30s).
        **kwargs: Keyword arguments for fn.
    """
    _name = task_name or getattr(fn, "__name__", "unknown")

    def _wrapper() -> None:
        start = time.time()
        logger.debug(f"[bg-task] START: {_name}")
        with _bg_task_metrics_lock:
            _bg_task_metrics["submitted"] += 1
        try:
            fn(*args, **kwargs)
            elapsed = (time.time() - start) * 1000
            logger.debug(f"[bg-task] DONE: {_name} ({elapsed:.0f}ms)")
            with _bg_task_metrics_lock:
                _bg_task_metrics["completed"] += 1
        except Exception as exc:
            elapsed = (time.time() - start) * 1000
            logger.error(
                f"[bg-task] FAILED: {_name} after {elapsed:.0f}ms: {exc}",
                exc_info=True,
            )
            with _bg_task_metrics_lock:
                _bg_task_metrics["failed"] += 1

    try:
        future = _background_executor.submit(_wrapper)

        # Add timeout callback -- cancels if still running after timeout
        def _check_timeout(f: Any) -> None:
            if not f.done():
                logger.warning(f"[bg-task] TIMEOUT: {_name} exceeded {timeout}s")
                with _bg_task_metrics_lock:
                    _bg_task_metrics["timed_out"] += 1

        # Schedule timeout check using a timer thread
        timer = threading.Timer(timeout, _check_timeout, args=(future,))
        timer.daemon = True
        timer.start()
    except RuntimeError:
        # Executor was shut down
        logger.warning(f"[bg-task] Executor shutdown, cannot submit: {_name}")


def get_background_task_metrics() -> dict[str, Any]:
    """Return background task execution metrics.

    Returns:
        Dict with submitted, completed, failed, timed_out counts.
    """
    with _bg_task_metrics_lock:
        return {
            **_bg_task_metrics,
            "executor_threads": _background_executor._max_workers,
        }


def _copilot_suggest(field: str, partial_input: str, context: dict) -> list[dict]:
    """Generate inline co-pilot suggestions for the media plan generator.

    Uses API integrations for data-backed suggestions and the LLM router
    for brief improvement tips. Results are cached for 60 seconds.

    Args:
        field: One of 'roles', 'locations', 'channels', 'brief'.
        partial_input: The user's current partial input text.
        context: Additional context (industry, budget, existing roles/locations).

    Returns:
        List of suggestion dicts: [{"text": "...", "confidence": 0.95, "reason": "..."}]
    """
    # Build cache key
    cache_key = (
        f"{field}:{partial_input}:{json.dumps(context, sort_keys=True, default=str)}"
    )
    cache_hash = hashlib.md5(cache_key.encode()).hexdigest()

    # Check cache
    with _copilot_cache_lock:
        if cache_hash in _copilot_cache:
            cached_suggestions, cached_at = _copilot_cache[cache_hash]
            if time.time() - cached_at < _COPILOT_CACHE_TTL:
                return cached_suggestions

    suggestions: list[dict] = []
    industry = context.get("industry") or ""
    budget = context.get("budget") or ""
    existing_roles = context.get("roles") or []
    existing_locations = context.get("locations") or []

    try:
        if field == "roles":
            suggestions = _copilot_suggest_roles(
                partial_input, industry, existing_roles
            )
        elif field == "locations":
            suggestions = _copilot_suggest_locations(
                partial_input, industry, budget, existing_roles, existing_locations
            )
        elif field == "channels":
            suggestions = _copilot_suggest_channels(
                industry, existing_roles, existing_locations, budget
            )
        elif field == "brief":
            suggestions = _copilot_suggest_brief(partial_input, context)
    except Exception as e:
        logger.error(
            "Copilot suggestion failed for field=%s: %s", field, e, exc_info=True
        )

    # Store in cache
    with _copilot_cache_lock:
        _copilot_cache[cache_hash] = (suggestions, time.time())
        # Evict stale + oldest entries if cache exceeds max
        if len(_copilot_cache) > _CACHE_MAX_ENTRIES:
            now = time.time()
            stale = [
                k
                for k, (_, t) in _copilot_cache.items()
                if now - t > _COPILOT_CACHE_TTL
            ]
            for k in stale:
                _copilot_cache.pop(k, None)
            # If still over limit, evict oldest
            if len(_copilot_cache) > _CACHE_MAX_ENTRIES:
                sorted_keys = sorted(_copilot_cache, key=lambda k: _copilot_cache[k][1])
                for k in sorted_keys[: len(_copilot_cache) - _CACHE_MAX_ENTRIES]:
                    _copilot_cache.pop(k, None)

    return suggestions


def _copilot_suggest_roles(
    partial: str, industry: str, existing: list[str]
) -> list[dict]:
    """Suggest roles based on industry and partial input using O*NET + Adzuna data."""
    suggestions: list[dict] = []

    # Try O*NET occupation search for data-backed role suggestions
    if _api_integrations_available and _api_onet and partial and len(partial) >= 2:
        try:
            onet_results = _api_onet.search_occupations(partial)
            if onet_results:
                for occ in onet_results[:5]:
                    title = occ.get("title") or occ.get("name") or ""
                    if title and title not in existing:
                        suggestions.append(
                            {
                                "text": title,
                                "confidence": 0.90,
                                "reason": "From O*NET occupation database",
                            }
                        )
        except Exception as e:
            logger.error("Copilot O*NET role lookup failed: %s", e, exc_info=True)

    # Try Adzuna for trending roles in the industry
    if (
        _api_integrations_available
        and _api_adzuna
        and industry
        and len(suggestions) < 3
    ):
        try:
            industry_label = industry.replace("_", " ").title()
            search_query = f"{partial} {industry_label}" if partial else industry_label
            adzuna_results = _api_adzuna.search_jobs(
                search_query, location="us", results_per_page=10
            )
            if adzuna_results and isinstance(adzuna_results, dict):
                seen_titles: set[str] = set()
                for job in adzuna_results.get("results") or []:
                    title = (job.get("title") or "").strip()
                    # Normalize: take first part before " - " or " at "
                    for sep in [" - ", " at ", " | ", ", "]:
                        if sep in title:
                            title = title.split(sep)[0].strip()
                    if (
                        title
                        and len(title) < 60
                        and title not in existing
                        and title.lower() not in seen_titles
                    ):
                        seen_titles.add(title.lower())
                        suggestions.append(
                            {
                                "text": title,
                                "confidence": 0.82,
                                "reason": f"Trending on job boards for {industry_label}",
                            }
                        )
                        if len(suggestions) >= 5:
                            break
        except Exception as e:
            logger.error("Copilot Adzuna role lookup failed: %s", e, exc_info=True)

    # Fallback: industry-based common roles
    if not suggestions and industry:
        _industry_roles: dict[str, list[str]] = {
            "healthcare_medical": [
                "Registered Nurse",
                "Medical Assistant",
                "Physician",
                "Pharmacist",
                "Physical Therapist",
            ],
            "technology": [
                "Software Engineer",
                "Data Scientist",
                "Product Manager",
                "DevOps Engineer",
                "UX Designer",
            ],
            "blue_collar_trades": [
                "Electrician",
                "Plumber",
                "HVAC Technician",
                "Welder",
                "CDL Truck Driver",
            ],
            "finance_banking": [
                "Financial Analyst",
                "Accountant",
                "Risk Manager",
                "Compliance Officer",
                "Loan Officer",
            ],
            "retail_hospitality": [
                "Store Manager",
                "Sales Associate",
                "Restaurant Manager",
                "Chef",
                "Front Desk Agent",
            ],
            "manufacturing": [
                "Production Supervisor",
                "Quality Inspector",
                "CNC Operator",
                "Maintenance Technician",
                "Supply Chain Analyst",
            ],
            "education": [
                "Teacher",
                "School Counselor",
                "Academic Advisor",
                "Curriculum Developer",
                "Special Education Teacher",
            ],
            "government_defense": [
                "Program Analyst",
                "Intelligence Analyst",
                "Cybersecurity Specialist",
                "Logistics Manager",
                "Contract Specialist",
            ],
        }
        role_list = _industry_roles.get(industry, [])
        for role in role_list:
            if role not in existing:
                suggestions.append(
                    {
                        "text": role,
                        "confidence": 0.75,
                        "reason": f"Common role in {industry.replace('_', ' ').title()}",
                    }
                )

    return suggestions[:5]


def _copilot_suggest_locations(
    partial: str,
    industry: str,
    budget: str,
    existing_roles: list[str],
    existing_locations: list[str],
) -> list[dict]:
    """Suggest locations based on roles, industry, and budget context."""
    suggestions: list[dict] = []

    # Use Adzuna to find locations with job postings for the roles
    if _api_integrations_available and _api_adzuna and existing_roles:
        try:
            role_query = existing_roles[0] if existing_roles else ""
            if role_query:
                adzuna_results = _api_adzuna.search_jobs(
                    role_query, location="us", results_per_page=20
                )
                if adzuna_results and isinstance(adzuna_results, dict):
                    location_counts: dict[str, int] = {}
                    for job in adzuna_results.get("results") or []:
                        loc = job.get("location") or {}
                        display = ""
                        area_list = loc.get("area") or []
                        if isinstance(area_list, list) and len(area_list) >= 2:
                            display = f"{area_list[-1]}, {area_list[0]}"
                        elif isinstance(area_list, list) and area_list:
                            display = area_list[0]
                        display_name = (loc.get("display_name") or display).strip()
                        if display_name and display_name not in existing_locations:
                            location_counts[display_name] = (
                                location_counts.get(display_name, 0) + 1
                            )
                    # Sort by frequency
                    for loc_name, count in sorted(
                        location_counts.items(), key=lambda x: -x[1]
                    )[:5]:
                        suggestions.append(
                            {
                                "text": loc_name,
                                "confidence": min(0.70 + count * 0.05, 0.99),
                                "reason": f"High demand for {existing_roles[0]} ({count} openings found)",
                            }
                        )
        except Exception as e:
            logger.error("Copilot Adzuna location lookup failed: %s", e, exc_info=True)

    # Fallback: popular hiring hubs
    if not suggestions:
        _top_hubs = [
            ("New York, NY", "Largest US talent market"),
            ("San Francisco, CA", "Tech talent hub, premium CPA"),
            ("Dallas, TX", "Fast-growing market, competitive CPA"),
            ("Chicago, IL", "Midwest talent hub, moderate cost"),
            ("London, UK", "Largest European talent market"),
        ]
        for city, reason in _top_hubs:
            if city not in existing_locations:
                suggestions.append({"text": city, "confidence": 0.70, "reason": reason})

    return suggestions[:5]


def _copilot_suggest_channels(
    industry: str, roles: list[str], locations: list[str], budget: str
) -> list[dict]:
    """Suggest channel categories based on industry, roles, and budget."""
    suggestions: list[dict] = []

    # Use benchmark_registry for data-backed channel recommendations
    try:
        registry_key = industry.replace("_", " ").lower() if industry else ""
        if registry_key:
            bench = benchmark_registry.get_channel_benchmark(
                "programmatic", registry_key
            )
            if bench:
                avg_cpa = bench.get("cpa") or bench.get("avg_cpa")
                if avg_cpa:
                    suggestions.append(
                        {
                            "text": "Programmatic/DSP Advertising",
                            "confidence": 0.88,
                            "reason": (
                                f"Avg CPA ${avg_cpa:.0f} for {industry.replace('_', ' ').title()}"
                                if isinstance(avg_cpa, (int, float))
                                else f"Recommended for {industry.replace('_', ' ').title()}"
                            ),
                        }
                    )
    except Exception as e:
        logger.error("Copilot benchmark lookup failed: %s", e, exc_info=True)

    # Industry-specific channel recommendations
    _channel_recs: dict[str, list[tuple[str, float, str]]] = {
        "healthcare_medical": [
            (
                "Niche & Industry-Specific Boards",
                0.92,
                "Health eCareers, Nurse.com deliver 40% lower CPA for healthcare",
            ),
            (
                "Social Media Channels",
                0.85,
                "Facebook Jobs effective for nursing/allied health recruitment",
            ),
        ],
        "technology": [
            (
                "Niche & Industry-Specific Boards",
                0.90,
                "Dice, Built In, Stack Overflow for targeted tech talent",
            ),
            (
                "Programmatic/DSP",
                0.88,
                "Programmatic reaches passive tech candidates at scale",
            ),
            (
                "Employer Branding",
                0.82,
                "Glassdoor/Indeed presence critical for tech hiring",
            ),
        ],
        "blue_collar_trades": [
            (
                "Regional & Local Job Boards",
                0.93,
                "Craigslist + local boards deliver best CPA for trades",
            ),
            (
                "Social Media Channels",
                0.87,
                "Facebook Jobs excels for blue collar recruitment",
            ),
        ],
        "retail_hospitality": [
            (
                "Regional & Local Job Boards",
                0.90,
                "Indeed + local boards for high-volume retail hiring",
            ),
            (
                "Social Media Channels",
                0.88,
                "Instagram/TikTok effective for Gen Z retail candidates",
            ),
        ],
    }
    for text, conf, reason in _channel_recs.get(industry, []):
        if not any(s["text"] == text for s in suggestions):
            suggestions.append({"text": text, "confidence": conf, "reason": reason})

    # If no industry-specific, give general recommendations
    if not suggestions:
        suggestions = [
            {
                "text": "Global Reach Job Boards",
                "confidence": 0.85,
                "reason": "Indeed + ZipRecruiter cover the widest candidate pool",
            },
            {
                "text": "Social Media Channels",
                "confidence": 0.80,
                "reason": "Social channels reach passive candidates effectively",
            },
            {
                "text": "Programmatic/DSP",
                "confidence": 0.78,
                "reason": "Automated ad placement optimizes CPA over time",
            },
        ]

    return suggestions[:5]


def _copilot_suggest_brief(partial: str, context: dict) -> list[dict]:
    """Suggest improvements to the campaign brief using LLM."""
    if not partial or len(partial.strip()) < 20:
        return [
            {
                "text": "Include the number of hires needed and timeline",
                "confidence": 0.90,
                "reason": "Specific targets enable precise budget allocation",
            },
            {
                "text": "Mention the hiring challenges (e.g., high turnover, niche skills)",
                "confidence": 0.85,
                "reason": "Challenges inform channel strategy and messaging",
            },
            {
                "text": "Add competitor names for competitive intelligence",
                "confidence": 0.80,
                "reason": "Competitor data enriches channel recommendations",
            },
        ]

    router = _lazy_llm_router()
    if not router:
        # Fallback: return generic tips based on what's missing from the brief
        tips = []
        lower = partial.lower()
        if not any(w in lower for w in ["hire", "recruit", "fill", "open"]):
            tips.append(
                {
                    "text": "Specify hiring volume (e.g., '200 hires in Q2')",
                    "confidence": 0.88,
                    "reason": "Hiring targets drive budget recommendations",
                }
            )
        if not any(w in lower for w in ["budget", "spend", "cost", "$"]):
            tips.append(
                {
                    "text": "Mention budget constraints or targets",
                    "confidence": 0.85,
                    "reason": "Budget context optimizes channel mix allocation",
                }
            )
        if not any(
            w in lower
            for w in ["challenge", "struggle", "issue", "problem", "turnover"]
        ):
            tips.append(
                {
                    "text": "Describe key hiring challenges",
                    "confidence": 0.82,
                    "reason": "Pain points inform strategy recommendations",
                }
            )
        return tips[:3] if tips else []

    try:
        industry = context.get("industry") or ""
        industry_label = industry.replace("_", " ")
        brief_excerpt = partial[:500]
        example_json = '[{"text": "Add timeline: Q2 2026", "confidence": 0.9, "reason": "Timeline enables seasonal optimization"}]'
        prompt = (
            f"You are a recruitment marketing strategist reviewing a campaign brief. "
            f"The user is writing a brief for a media plan in the {industry_label} industry.\n\n"
            f'Current brief:\n"{brief_excerpt}"\n\n'
            f"Suggest 2-3 specific improvements to make this brief more actionable for media planning. "
            f"Each suggestion should be a short phrase (under 60 chars) that the user can click to add context. "
            f"Return ONLY a JSON array of objects with keys: text, confidence (0-1), reason (one sentence).\n"
            f"Example: {example_json}"
        )
        # S50: TASK_PLAN_STRUCTURED for JSON output (Gemini best at structured),
        # 10s timeout to avoid blocking brief typing experience.
        result = router.call_llm(
            messages=[{"role": "user", "content": prompt}],
            system_prompt="Return only valid JSON array. No markdown, no explanation.",
            task_type=getattr(router, "TASK_PLAN_STRUCTURED", "plan_structured"),
            max_tokens=300,
            timeout_budget=10.0,
        )
        response_text = result.get("text") or ""
        # Parse JSON from response
        # Strip markdown code fences if present
        cleaned = response_text.strip()
        if cleaned.startswith("```"):
            cleaned = cleaned.split("\n", 1)[-1] if "\n" in cleaned else cleaned[3:]
            if cleaned.endswith("```"):
                cleaned = cleaned[:-3]
            cleaned = cleaned.strip()
        parsed = json.loads(cleaned)
        if isinstance(parsed, list):
            suggestions = []
            for item in parsed[:3]:
                if isinstance(item, dict) and item.get("text"):
                    suggestions.append(
                        {
                            "text": str(item["text"])[:80],
                            "confidence": float(item.get("confidence") or 0.80),
                            "reason": str(item.get("reason") or "")[:120],
                        }
                    )
            return suggestions
    except json.JSONDecodeError:
        logger.warning("Copilot brief suggestion: LLM returned non-JSON response")
    except Exception as e:
        logger.error("Copilot brief suggestion LLM call failed: %s", e, exc_info=True)

    return []


# ═══════════════════════════════════════════════════════════════════════════════
# PLATFORM FRAGMENT CACHE (HTML fragments for SPA shell)
# ═══════════════════════════════════════════════════════════════════════════════
_fragment_cache: dict[str, str] = {}
_fragment_cache_lock = threading.Lock()

# Map fragment names to template files and extraction selectors
_FRAGMENT_MAP: dict[str, str] = {
    # ── Core module fragments (10 modules) ──
    "media-plan": "index.html",
    "quick-plan": "quick-plan.html",
    "quick-brief": "quick-brief.html",
    "social-plan": "social-plan.html",
    "creative-ai": "creative-ai.html",
    "ab-testing": "ab-testing.html",
    "competitive": "competitive.html",
    "market-intel": "market-intel.html",
    "market-pulse": "market-pulse.html",
    "vendor-iq": "vendor-iq.html",
    "simulator": "simulator.html",
    "roi-calculator": "roi-calculator.html",
    "tracker": "tracker.html",
    "post-campaign": "post-campaign.html",
    "talent-heatmap": "talent-heatmap.html",
    "hire-signal": "hire-signal.html",
    "payscale-sync": "payscale-sync.html",
    "skill-target": "skill-target.html",
    "applyflow": "applyflow-demo.html",
    "compliance-guard": "compliance-guard.html",
    "audit": "audit.html",
    "nova": "nova.html",
    # ── SlotOps ──
    "slotops": "slotops.html",
    # ── Utility fragments ──
    "api-portal": "api-portal.html",
    "dashboard": "hub.html",
    # ── Backward-compat alias ──
    "applyflow-demo": "applyflow-demo.html",
}
_INSIGHTS_LLM_TIMEOUT = 10  # seconds
_insights_executor = ThreadPoolExecutor(
    max_workers=2, thread_name_prefix="llm-insights"
)


def _generate_product_insights(
    product_name: str,
    result_summary: dict,
    context: str = "",
    max_tokens: int = 250,
) -> str:
    """Generate optional AI insights for a product endpoint.

    Uses the LLM router with caching (1h TTL) and a 10-second timeout.
    Returns empty string on any failure -- never blocks the main response.

    Args:
        product_name: Human-readable product name (e.g. "Budget Simulator").
        result_summary: Dict of key result data to base insights on.
        context: Additional context for the prompt.
        max_tokens: Max tokens for LLM response.

    Returns:
        AI-generated insights string, or empty string on failure.
    """
    router = _lazy_llm_router()
    if not router:
        return ""

    # Build cache key from product name + truncated summary hash
    summary_str = json.dumps(result_summary, sort_keys=True, default=str)[:1000]
    cache_key = hashlib.md5(
        f"{product_name}:{summary_str}".encode(), usedforsecurity=False
    ).hexdigest()

    # Check cache
    with _insights_cache_lock:
        cached = _insights_cache.get(cache_key)
        if cached:
            text, ts = cached
            if time.time() - ts < _INSIGHTS_CACHE_TTL:
                return text
            else:
                del _insights_cache[cache_key]

    def _call_llm_sync() -> str:
        """Run LLM call in thread for timeout control."""
        # S48: Route insights to TASK_INTELLIGENCE_SUMMARY (Gemini Flash Lite)
        # instead of TASK_STRUCTURED (full Gemini). Insights are short summaries
        # that don't need the full model.
        task_type = getattr(router, "TASK_INTELLIGENCE_SUMMARY", "intelligence_summary")

        # ── RAG: Retrieve relevant knowledge base context ──
        rag_context = ""
        if _vector_search_available and _vector_search:
            try:
                # Build search query from product data
                rag_query_parts: list[str] = []
                for key in ("job_title", "industry", "location", "roles", "channel"):
                    val = result_summary.get(key)
                    if val:
                        rag_query_parts.append(
                            str(val)
                            if not isinstance(val, list)
                            else " ".join(str(v) for v in val)
                        )
                if not rag_query_parts and context:
                    rag_query_parts.append(context[:200])
                if not rag_query_parts:
                    rag_query_parts.append(product_name)
                rag_query = " ".join(rag_query_parts)[:300]
                kb_results = _vector_search(rag_query, top_k=3)
                if kb_results:
                    snippets = [
                        r.get("text") or ""
                        for r in kb_results
                        if (r.get("score") or 0) > 0.3
                    ]
                    if snippets:
                        rag_context = (
                            "Relevant knowledge base context:\n"
                            + "\n---\n".join(snippets[:3])
                            + "\n\n"
                        )
            except Exception as rag_err:
                logger.debug(f"RAG retrieval for {product_name} skipped: {rag_err}")

        prompt = f"You are analyzing {product_name} output for a recruitment media planner.\n\n"
        if rag_context:
            prompt += rag_context
        prompt += f"Data summary:\n{json.dumps(result_summary, indent=2, default=str)[:1500]}\n\n"
        if context:
            prompt += f"Context: {context}\n\n"
        prompt += (
            "Provide exactly 3 brief, actionable insights as a numbered list. "
            "Each insight should be 1-2 sentences. Focus on practical recommendations "
            "a recruitment media planner can act on immediately."
        )
        # S50: 10s timeout_budget matches _INSIGHTS_LLM_TIMEOUT thread timeout
        result = router.call_llm(
            messages=[{"role": "user", "content": prompt}],
            system_prompt=(
                "You are a senior recruitment marketing strategist. "
                "Give concise, data-driven insights. No fluff, no preamble."
            ),
            task_type=task_type,
            max_tokens=max_tokens,
            timeout_budget=10.0,
        )
        return result.get("text") or ""

    try:
        future = _insights_executor.submit(_call_llm_sync)
        text = future.result(timeout=_INSIGHTS_LLM_TIMEOUT)
        # Store in cache
        if text:
            with _insights_cache_lock:
                _insights_cache[cache_key] = (text, time.time())
                # Evict stale + oldest entries if cache exceeds max
                if len(_insights_cache) > _CACHE_MAX_ENTRIES:
                    now = time.time()
                    stale = [
                        k
                        for k, (_, ts) in _insights_cache.items()
                        if now - ts > _INSIGHTS_CACHE_TTL
                    ]
                    for k in stale:
                        del _insights_cache[k]
                    # If still over limit, evict oldest
                    if len(_insights_cache) > _CACHE_MAX_ENTRIES:
                        sorted_keys = sorted(
                            _insights_cache, key=lambda k: _insights_cache[k][1]
                        )
                        for k in sorted_keys[
                            : len(_insights_cache) - _CACHE_MAX_ENTRIES
                        ]:
                            del _insights_cache[k]
        return text
    except TimeoutError:
        logger.warning(
            f"LLM insights timed out for {product_name} (>{_INSIGHTS_LLM_TIMEOUT}s)"
        )
        return ""
    except Exception as exc:
        logger.error(f"LLM insights failed for {product_name}: {exc}", exc_info=True)
        return ""


# ═══════════════════════════════════════════════════════════════════════════════
# SENTRY ERROR TRACKING
# ═══════════════════════════════════════════════════════════════════════════════
_SENTRY_DSN = (os.environ.get("SENTRY_DSN") or "").strip()
if _SENTRY_DSN:
    try:
        import sentry_sdk

        def _sentry_before_send(event: dict, hint: dict) -> dict | None:
            """Filter out noisy/expected errors from Sentry.

            Drops events that are operational noise rather than real bugs:
            - Auth/credential errors (401/403) for external APIs
            - FRED HTTP 400 errors -- bad series IDs are data issues
            - Rate limiting (429) -- transient, handled by retry logic
            - Timeouts during startup/warmup -- expected on cold starts
            - Missing API keys -- config issues, not bugs
            - Transient network errors (connection reset, DNS)
            - API integration test failures -- expected when keys missing
            - Startup failures during warmup period
            - S63: exception-class filter (TimeoutError, ConnectionError, etc.)
              -- existing string matching missed most of these because they
              arrive as exception events, not message events.
            """
            # S63: Exception-class filter. Hit BEFORE string matching so we
            # catch events that have no meaningful .message (most of the 4,179
            # errors consuming the Joveo Sentry quota).
            exc_info = hint.get("exc_info") if hint else None
            if exc_info and len(exc_info) >= 2:
                exc_type = exc_info[0]
                exc_value = exc_info[1]
                exc_str = str(exc_value or "").lower()

                # Drop all transient network/timeout exceptions outright --
                # these are handled by retry/circuit breakers and are not bugs.
                import socket as _socket

                if (
                    exc_type
                    and isinstance(exc_type, type)
                    and issubclass(
                        exc_type,
                        (TimeoutError, ConnectionError, _socket.timeout),
                    )
                ):
                    return None

                # Drop urllib HTTPError 4xx (client/config issues, not bugs).
                exc_name = exc_type.__name__ if exc_type else ""
                if exc_name == "HTTPError" and any(
                    code in exc_str for code in ("400", "401", "403", "404", "429")
                ):
                    return None

                # Drop flaky JSON decode errors from upstream APIs returning HTML.
                if exc_name in ("JSONDecodeError", "ValueError") and any(
                    kw in exc_str
                    for kw in ("expecting value", "no json", "<html", "<!doctype")
                ):
                    return None

            # S63: Logger-name filter. These modules have retry + fallback
            # semantics; their errors are operational, not bugs. The log line
            # still appears in the Sentry breadcrumb trail attached to real
            # events, so signal is preserved for debugging.
            _noisy_loggers = {
                "llm_router",
                "data_matrix_monitor",
                "auto_qc",
                "data_enrichment",
                "vector_search",
                "firecrawl_enrichment",
                "tavily_search",
                "calendar_sync",
                "market_signals",
                "nova_proactive",
            }
            logger_name = (event.get("logger") or "").split(".")[0].lower()
            if logger_name in _noisy_loggers and event.get("level") in (
                "error",
                "warning",
            ):
                return None

            message = (event.get("message") or "").lower()
            logentry = ((event.get("logentry") or {}).get("message") or "").lower()
            combined = f"{message} {logentry}"

            # 1. Drop 401/403 credential errors -- config issues, not bugs
            if "http 401" in combined or "http 403" in combined:
                if any(
                    kw in combined
                    for kw in (
                        "onetcenter",
                        "auth/credential",
                        "rate-limited",
                        "api key",
                        "credential",
                        "unauthorized",
                    )
                ):
                    return None

            # 2. Drop FRED HTTP 400 errors -- bad/unknown series IDs are data issues, not bugs
            if "http 400" in combined and "stlouisfed" in combined:
                return None

            # 3. Drop 429 rate-limit errors -- transient, handled by retry
            if (
                "http 429" in combined
                or "rate-limited" in combined
                or "rate limit" in combined
            ):
                return None

            # 4. Drop transient network errors -- not bugs
            _transient_patterns = (
                "timeout fetching",
                "timeout posting",
                "os error fetching",
                "os error posting",
                "connection reset",
                "connection refused",
                "name resolution",
                "ssl verification failed",
                "timed out",
                "[errno 104]",
                "[errno 111]",
                "urlopen error",
                "temporary failure",
            )
            if any(pat in combined for pat in _transient_patterns):
                return None

            # 5. Drop API integration test failures -- expected when keys missing
            _api_test_patterns = (
                "fred test failed",
                "adzuna test failed",
                "jooble test failed",
                "o*net test failed",
                "bea test failed",
                "census test failed",
                "usajobs test failed",
                "bls test failed",
            )
            if any(pat in combined for pat in _api_test_patterns):
                return None

            # 6. Drop startup/warmup noise
            _startup_patterns = (
                "knowledge base pre-warm failed",
                "vector index build failed",
                "data refresh failed",
                "feature store init failed",
                "datamatrixmonitor: check failed",
                "api_keys probe failed",
                "enrichment cycle crashed",
            )
            if any(pat in combined for pat in _startup_patterns):
                if not _DEPLOY_WARMUP_COMPLETE:
                    return None

            # 7. Drop missing API key warnings logged as errors
            if "api key" in combined and (
                "not set" in combined or "missing" in combined
            ):
                return None

            return event

        sentry_sdk.init(
            dsn=_SENTRY_DSN,
            traces_sample_rate=0.1,  # 10% of requests get traces
            profiles_sample_rate=0.1,  # 10% get profiling
            environment="production" if os.environ.get("RENDER") else "local",
            release="media-plan-generator@4.0.0",
            send_default_pii=False,  # Never send PII
            attach_stacktrace=True,
            before_send=_sentry_before_send,
        )
        logger.info("Sentry error tracking initialized (SDK %s)", sentry_sdk.VERSION)
        _SENTRY_AVAILABLE = True
        # NOTE: Removed startup capture_message -- it creates noise in Sentry
        # Sentry will auto-detect the project on the first real error event
    except ImportError:
        logger.warning("sentry-sdk not installed -- error tracking disabled")
        _SENTRY_AVAILABLE = False
    except Exception as _sentry_err:
        logger.warning("Sentry init failed: %s", _sentry_err)
        _SENTRY_AVAILABLE = False
else:
    _SENTRY_AVAILABLE = False


# ── Sentry Performance Helpers (SDK 2.x API) ──────────────────────────────
from contextlib import nullcontext as _nullctx


def _sentry_start_txn(op: str, name: str):
    """Start a Sentry root span (transaction) using SDK 2.x API.

    In SDK 2.x, ``sentry_sdk.start_span()`` creates a root span (transaction)
    when no parent span exists on the current scope. Child spans are created
    the same way -- the SDK auto-parents them.

    Returns (root_span_cm, child_span_fn).
    """
    if not _SENTRY_AVAILABLE:
        return _nullctx(), lambda op, name: _nullctx()
    try:
        root_span = sentry_sdk.start_span(op=op, name=name)

        def _child(child_op: str, child_name: str):
            return sentry_sdk.start_span(op=child_op, name=child_name)

        return root_span, _child
    except Exception:
        return _nullctx(), lambda op, name: _nullctx()


# ═══════════════════════════════════════════════════════════════════════════════
# KNOWLEDGE BASE LOADER
# ═══════════════════════════════════════════════════════════════════════════════

_knowledge_base = None

# Prefer kb_loader's thread-safe implementation (uses _kb_lock).
# Falls back to the inline definition below if kb_loader is unavailable.
try:
    from kb_loader import load_knowledge_base  # thread-safe, has _kb_lock
except ImportError:
    pass  # inline load_knowledge_base() defined below serves as fallback

# Mapping from legacy industry keys (used in INDUSTRY_NAICS_MAP) to KB keys
# in data/recruitment_industry_knowledge.json -> industry_specific_benchmarks
_INDUSTRY_KEY_TO_KB_KEY = {
    "healthcare": "healthcare",
    "healthcare_medical": "healthcare",
    "mental_health": "healthcare",
    "technology": "technology",
    "tech_engineering": "technology",
    "telecom": "technology",
    "telecommunications": "technology",
    "retail": "retail_hospitality",
    "retail_consumer": "retail_hospitality",
    "hospitality": "retail_hospitality",
    "hospitality_travel": "retail_hospitality",
    "food_beverage": "retail_hospitality",
    "construction": "construction_infrastructure",
    "construction_real_estate": "construction_infrastructure",
    "transportation": "transportation_logistics",
    "logistics_supply_chain": "transportation_logistics",
    "maritime": "transportation_logistics",
    "maritime_marine": "transportation_logistics",
    "rideshare": "transportation_logistics",
    "manufacturing": "manufacturing",
    "automotive": "manufacturing",
    "blue_collar": "transportation_logistics",
    "blue_collar_trades": "transportation_logistics",
    "aerospace": "manufacturing",
    "aerospace_defense": "manufacturing",
    "pharma": "healthcare",
    "pharma_biotech": "healthcare",
    "finance": "financial_services",
    "finance_banking": "financial_services",
    "insurance": "financial_services",
    "government": "government_utilities",
    "military_recruitment": "government_utilities",
    "energy": "government_utilities",
    "energy_utilities": "government_utilities",
    "education": "government_utilities",
    "professional_services": "technology",
    "legal_services": "financial_services",
    "nonprofit": "government_utilities",
    "general": "retail_hospitality",
    "general_entry_level": "retail_hospitality",
    "media_entertainment": "technology",
}


def load_knowledge_base() -> dict:
    """Load and merge all knowledge base files into unified dict.

    Loads 10 JSON data files from the ``data/`` directory, each into its own
    section key.  The ``core`` section (recruitment_industry_knowledge.json) is
    additionally merged into the top level for backward compatibility so that
    existing code referencing ``kb["benchmarks"]``, ``kb["salary_trends"]``,
    etc. continues to work.

    Uses a module-level cache so the files are read at most once.

    Returns:
        Merged dict with section keys + backward-compat top-level keys,
        or a minimal dict on failure.
    """
    global _knowledge_base
    if _knowledge_base is not None:
        return _knowledge_base

    files = {
        "core": "recruitment_industry_knowledge.json",
        "platform_intelligence": "platform_intelligence_deep.json",
        "recruitment_benchmarks": "recruitment_benchmarks_deep.json",
        "recruitment_strategy": "recruitment_strategy_intelligence.json",
        "regional_hiring": "regional_hiring_intelligence.json",
        "supply_ecosystem": "supply_ecosystem_intelligence.json",
        "workforce_trends": "workforce_trends_intelligence.json",
        "white_papers": "industry_white_papers.json",
        "joveo_2026_benchmarks": "joveo_2026_benchmarks.json",
        "google_ads_benchmarks": "google_ads_2025_benchmarks.json",
        "external_benchmarks": "external_benchmarks_2025.json",
        "client_media_plans": "client_media_plans_kb.json",
        "linkedin_performance": "linkedin_performance_benchmarks.json",
        "linkedin_industry_benchmarks": "linkedin_industry_benchmarks.json",
        "craigslist_performance": "craigslist_performance_benchmarks.json",
        "international_benchmarks": "international_benchmarks_2026.json",
        "slotops_benchmarks": "slotops_benchmarks_summary.json",
        # S50: 15 previously unindexed data files
        "adzuna_benchmarks": "adzuna_benchmarks.json",
        "channel_benchmarks_live": "channel_benchmarks_live.json",
        "channels_db": "channels_db.json",
        "competitor_careers": "competitor_careers.json",
        "fred_indicators": "fred_indicators.json",
        "google_trends": "google_trends.json",
        "h1b_salary_intelligence": "h1b_salary_intelligence.json",
        "job_density_metros": "job_density_metros.json",
        "job_posting_volumes": "job_posting_volumes.json",
        "joveo_publishers": "joveo_publishers.json",
        "live_market_data": "live_market_data.json",
        "market_trends_live": "market_trends_live.json",
        "platform_ad_specs": "platform_ad_specs.json",
        "seasonal_hiring_trends": "seasonal_hiring_trends.json",
        "global_supply": "global_supply.json",
    }

    kb = {}
    data_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "data")
    loaded_count = 0
    for section_key, filename in files.items():
        fpath = os.path.join(data_dir, filename)
        try:
            with open(fpath, "r", encoding="utf-8") as f:
                kb[section_key] = json.load(f)
                loaded_count += 1
                logger.info("KB loaded %s (%s)", section_key, filename)
        except FileNotFoundError:
            kb[section_key] = {}
            logger.warning("KB file not found: %s", filename)
        except json.JSONDecodeError as e:
            kb[section_key] = {}
            logger.error("KB JSON error in %s: %s", filename, e)
        except Exception as e:
            kb[section_key] = {}
            logger.error("KB load error for %s: %s", filename, e)

    # Backward compatibility: merge core keys to top level so existing
    # code that accesses kb["benchmarks"], kb["salary_trends"], etc. still works
    core = kb.get("core", {})
    for k, v in core.items():
        if k not in kb:  # don't overwrite section keys
            kb[k] = v

    # ── Data Freshness Validation ──
    # Check last_updated metadata in each KB section and warn if data is
    # older than 90 days.  This prevents silently serving stale benchmarks
    # when JSON files haven't been regenerated in a long time.
    stale_sections = []
    try:
        today = datetime.datetime.now()
        max_age_days = 90
        for section_key, section_data in kb.items():
            if not isinstance(section_data, dict):
                continue
            # Check top-level and nested metadata for last_updated
            last_updated_str = None
            if isinstance(section_data.get("metadata"), dict):
                last_updated_str = section_data["metadata"].get("last_updated")
            if not last_updated_str:
                last_updated_str = section_data.get("last_updated")
            if last_updated_str and isinstance(last_updated_str, str):
                try:
                    lu_date = datetime.datetime.strptime(
                        last_updated_str[:10], "%Y-%m-%d"
                    )
                    age_days = (today - lu_date).days
                    if age_days > max_age_days:
                        stale_sections.append((section_key, last_updated_str, age_days))
                except (ValueError, TypeError):
                    pass
        if stale_sections:
            for skey, sdate, sage in stale_sections:
                logger.warning(
                    "KB DATA FRESHNESS WARNING: '%s' last updated %s "
                    "(%d days ago, threshold=%d days)",
                    skey,
                    sdate,
                    sage,
                    max_age_days,
                )
            kb["_freshness_warnings"] = [
                {"section": s, "last_updated": d, "age_days": a}
                for s, d, a in stale_sections
            ]
    except Exception as e:
        logger.warning("KB freshness check failed (non-fatal): %s", e)

    # ── Load Joveo category spend data (CSV) ──
    try:
        import csv

        _joveo_spends_path = os.path.join(
            data_dir, "joveo_category_spends_oct25_feb26.csv"
        )
        if os.path.exists(_joveo_spends_path):
            _joveo_spends: list[dict[str, str]] = []
            with open(_joveo_spends_path, "r", encoding="utf-8-sig") as f:
                reader = csv.DictReader(f)
                for row in reader:
                    _joveo_spends.append(row)
            kb["joveo_category_spends"] = {
                "data": _joveo_spends,
                "record_count": len(_joveo_spends),
                "description": "Joveo category spend data Oct 2025 - Feb 2026 (433 subcategories)",
                "columns": list(_joveo_spends[0].keys()) if _joveo_spends else [],
            }
            logger.info("Loaded Joveo category spends: %d records", len(_joveo_spends))
    except FileNotFoundError:
        logger.warning("Joveo category spends CSV not found")
    except Exception as e:
        logger.error("Failed to load Joveo category spends: %s", e, exc_info=True)

    # ── KB Memory Usage Tracking (#14) ──
    try:
        kb_json_bytes = len(json.dumps(kb).encode("utf-8"))
        logger.info(
            "Knowledge base loaded: %d/%d files, %d total keys, ~%.1f MB in memory",
            loaded_count,
            len(files),
            len(kb),
            kb_json_bytes / 1_048_576,
        )
        if kb_json_bytes > 50 * 1_048_576:  # warn above 50 MB
            logger.warning(
                "KB memory usage HIGH: %.1f MB — consider lazy loading",
                kb_json_bytes / 1_048_576,
            )
    except Exception as mem_err:
        logger.info(
            "Knowledge base loaded: %d/%d files, %d total keys (memory tracking failed: %s)",
            loaded_count,
            len(files),
            len(kb),
            mem_err,
        )

    # ── KB Data Quality Validation (#16) ──
    try:
        quality_issues: list[dict] = []
        for qk, qv in kb.items():
            if qk.startswith("_"):
                continue
            issues: list[str] = []
            if qv is None:
                issues.append("null_data")
            elif isinstance(qv, dict) and len(qv) == 0:
                issues.append("empty_dict")
            elif isinstance(qv, list) and len(qv) == 0:
                issues.append("empty_list")
            if issues:
                quality_issues.append(
                    {"key": qk, "issues": issues, "type": type(qv).__name__}
                )
        if quality_issues:
            logger.warning(
                "KB quality issues in %d sections: %s",
                len(quality_issues),
                ", ".join(
                    f"{qi['key']}({','.join(qi['issues'])})" for qi in quality_issues
                ),
            )
            kb["_quality_issues"] = quality_issues
        else:
            logger.info("KB data quality check passed: all sections have data")
    except Exception as qe:
        logger.debug("KB quality check failed (non-fatal): %s", qe)

    _knowledge_base = kb
    return kb


# ═══════════════════════════════════════════════════════════════════════════════
# NAICS-BASED INDUSTRY CLASSIFICATION ENGINE
# ═══════════════════════════════════════════════════════════════════════════════
INDUSTRY_NAICS_MAP = {
    # Technology / Software (NAICS 51)
    "technology": {
        "naics": "51",
        "sector": "Technology & Software",
        "bls_sector": "Information",
        "talent_profile": "Professional / Technical Workforce",
        "legacy_key": "tech_engineering",
        "keywords": [
            "tech",
            "software",
            "saas",
            "cloud",
            "ai",
            "data",
            "cyber",
            "digital",
            "it services",
            "information technology",
        ],
    },
    # Healthcare (NAICS 62)
    "healthcare": {
        "naics": "62",
        "sector": "Healthcare & Life Sciences",
        "bls_sector": "Health Care",
        "talent_profile": "Clinical & Allied Health Workforce",
        "legacy_key": "healthcare_medical",
        "keywords": [
            "health",
            "hospital",
            "medical",
            "pharma",
            "biotech",
            "clinical",
            "nursing",
            "patient care",
            "life sciences",
        ],
    },
    # Finance (NAICS 52)
    "finance": {
        "naics": "52",
        "sector": "Financial Services & Insurance",
        "bls_sector": "Financial Activities",
        "talent_profile": "Professional / Financial Workforce",
        "legacy_key": "finance_banking",
        "keywords": [
            "bank",
            "finance",
            "insurance",
            "investment",
            "capital",
            "wealth",
            "fintech",
            "credit",
            "mortgage",
            "securities",
        ],
    },
    # Manufacturing (NAICS 31-33)
    "manufacturing": {
        "naics": "31-33",
        "sector": "Manufacturing & Industrial",
        "bls_sector": "Manufacturing",
        "talent_profile": "Skilled Trades & Production Workforce",
        "legacy_key": "automotive",
        "keywords": ["manufactur", "industrial", "production", "assembly", "factory"],
    },
    # Retail / E-commerce (NAICS 44-45)
    "retail": {
        "naics": "44-45",
        "sector": "Retail & E-Commerce",
        "bls_sector": "Retail Trade",
        "talent_profile": "General / Entry-Level & Hourly Workforce",
        "legacy_key": "retail_consumer",
        "keywords": [
            "retail",
            "e-commerce",
            "ecommerce",
            "store",
            "shop",
            "consumer",
            "merchandise",
            "grocery",
        ],
    },
    # Transportation / Logistics (NAICS 48-49)
    "transportation": {
        "naics": "48-49",
        "sector": "Transportation & Logistics",
        "bls_sector": "Transportation and Warehousing",
        "talent_profile": "Operations & Logistics Workforce",
        "legacy_key": "logistics_supply_chain",
        "keywords": [
            "transport",
            "logistics",
            "shipping",
            "delivery",
            "freight",
            "trucking",
            "courier",
            "warehouse",
            "supply chain",
            "fleet",
        ],
    },
    # Hospitality (NAICS 72)
    "hospitality": {
        "naics": "72",
        "sector": "Hospitality & Tourism",
        "bls_sector": "Leisure and Hospitality",
        "talent_profile": "Service & Hospitality Workforce",
        "legacy_key": "hospitality_travel",
        "keywords": [
            "hotel",
            "hospitality",
            "restaurant",
            "food service",
            "tourism",
            "travel",
            "lodging",
            "resort",
            "dining",
        ],
    },
    # Education (NAICS 61)
    "education": {
        "naics": "61",
        "sector": "Education & Training",
        "bls_sector": "Education",
        "talent_profile": "Academic & Educational Workforce",
        "legacy_key": "education",
        "keywords": [
            "education",
            "school",
            "university",
            "college",
            "academic",
            "teaching",
            "k-12",
            "k12",
            "training",
            "teacher",
            "principal",
            "isd",
        ],
    },
    # Professional Services (NAICS 54)
    "professional_services": {
        "naics": "54",
        "sector": "Professional & Business Services",
        "bls_sector": "Professional and Business Services",
        "talent_profile": "Professional / Consulting Workforce",
        "legacy_key": "legal_services",
        "keywords": [
            "consulting",
            "professional service",
            "accounting",
            "legal",
            "audit",
            "advisory",
            "staffing",
            "management consulting",
        ],
    },
    # Construction (NAICS 23)
    "construction": {
        "naics": "23",
        "sector": "Construction & Real Estate",
        "bls_sector": "Construction",
        "talent_profile": "Skilled Trades & Construction Workforce",
        "legacy_key": "construction_real_estate",
        "keywords": [
            "construction",
            "building",
            "real estate",
            "property",
            "architecture",
            "engineering firm",
            "contractor",
        ],
    },
    # Energy (NAICS 21, 22)
    "energy": {
        "naics": "21-22",
        "sector": "Energy & Utilities",
        "bls_sector": "Mining and Logging",
        "talent_profile": "Technical & Field Operations Workforce",
        "legacy_key": "energy_utilities",
        "keywords": [
            "energy",
            "oil",
            "gas",
            "solar",
            "renewable",
            "electric",
            "utility",
            "power",
            "mining",
            "petroleum",
        ],
    },
    # Government (NAICS 92)
    "government": {
        "naics": "92",
        "sector": "Government & Public Administration",
        "bls_sector": "Government",
        "talent_profile": "Public Sector Workforce",
        "legacy_key": "military_recruitment",
        "keywords": [
            "government",
            "federal",
            "state agency",
            "municipal",
            "public sector",
            "military",
            "defense",
            "civil service",
            "army",
            "navy",
            "air force",
            "marines",
        ],
    },
    # Telecommunications (NAICS 51 sub)
    "telecom": {
        "naics": "51",
        "sector": "Telecommunications",
        "bls_sector": "Information",
        "talent_profile": "Technical & Engineering Workforce",
        "legacy_key": "telecommunications",
        "keywords": [
            "telecom",
            "wireless",
            "broadband",
            "5g",
            "network",
            "communications",
        ],
    },
    # Agriculture (NAICS 11)
    "agriculture": {
        "naics": "11",
        "sector": "Agriculture & Food Production",
        "bls_sector": "Agriculture",
        "talent_profile": "Agricultural & Seasonal Workforce",
        "legacy_key": "food_beverage",
        "keywords": [
            "agriculture",
            "farming",
            "food production",
            "agri",
            "crop",
            "livestock",
        ],
    },
    # Nonprofit (NAICS 81)
    "nonprofit": {
        "naics": "81",
        "sector": "Nonprofit & Social Services",
        "bls_sector": "Other Services",
        "talent_profile": "Mission-Driven Workforce",
        "legacy_key": "general_entry_level",
        "keywords": [
            "nonprofit",
            "non-profit",
            "ngo",
            "charity",
            "foundation",
            "social services",
        ],
    },
    # General / Entry-Level (catch-all)
    "general": {
        "naics": "00",
        "sector": "General / Multi-Industry",
        "bls_sector": "Total Nonfarm",
        "talent_profile": "General / Entry-Level & Hourly Workforce",
        "legacy_key": "general_entry_level",
        "keywords": [
            "general",
            "entry level",
            "entry-level",
            "hourly",
            "temp",
            "seasonal",
            "staffing agency",
        ],
    },
    # Rideshare / Gig Economy (NAICS 48 sub)
    "rideshare": {
        "naics": "48",
        "sector": "Rideshare & Gig Economy",
        "bls_sector": "Transportation and Warehousing",
        "talent_profile": "Gig & Independent Contractor Workforce",
        "legacy_key": "logistics_supply_chain",
        "keywords": [
            "rideshare",
            "ride-share",
            "gig",
            "uber",
            "lyft",
            "doordash",
            "instacart",
            "platform",
            "on-demand",
        ],
    },
    # Blue Collar / Skilled Trades (NAICS 23/31-33 crossover)
    "blue_collar": {
        "naics": "31-33",
        "sector": "Blue Collar / Skilled Trades",
        "bls_sector": "Manufacturing",
        "talent_profile": "Skilled Trades & Production Workforce",
        "legacy_key": "blue_collar_trades",
        "keywords": [
            "blue collar",
            "skilled trade",
            "welder",
            "electrician",
            "plumber",
            "mechanic",
            "hvac",
            "trade",
        ],
    },
    # Maritime / Marine (NAICS 48 sub)
    "maritime": {
        "naics": "48",
        "sector": "Maritime & Marine",
        "bls_sector": "Transportation and Warehousing",
        "talent_profile": "Maritime & Marine Workforce",
        "legacy_key": "maritime_marine",
        "keywords": [
            "maritime",
            "marine",
            "shipping",
            "naval",
            "shipyard",
            "port",
            "offshore",
            "vessel",
            "seafar",
        ],
    },
    # Aerospace & Defense (NAICS 33 sub)
    "aerospace": {
        "naics": "33",
        "sector": "Aerospace & Defense",
        "bls_sector": "Manufacturing",
        "talent_profile": "Aerospace & Defense Workforce",
        "legacy_key": "aerospace_defense",
        "keywords": [
            "aerospace",
            "defense",
            "aviation",
            "aircraft",
            "missile",
            "satellite",
            "space",
        ],
    },
    # Pharma & Biotech (NAICS 32 sub)
    "pharma": {
        "naics": "32",
        "sector": "Pharma & Biotech",
        "bls_sector": "Manufacturing",
        "talent_profile": "Pharmaceutical & Research Workforce",
        "legacy_key": "pharma_biotech",
        "keywords": [
            "pharma",
            "biotech",
            "drug",
            "vaccine",
            "clinical trial",
            "biolog",
            "therapeutic",
        ],
    },
    # Insurance (NAICS 52 sub)
    "insurance": {
        "naics": "52",
        "sector": "Insurance",
        "bls_sector": "Financial Activities",
        "talent_profile": "Professional / Insurance Workforce",
        "legacy_key": "insurance",
        "keywords": ["insurance", "underwriter", "actuar", "claims", "policy"],
    },
    # Mental Health (NAICS 62 sub)
    "mental_health": {
        "naics": "62",
        "sector": "Mental Health & Behavioral",
        "bls_sector": "Health Care",
        "talent_profile": "Behavioral & Mental Health Workforce",
        "legacy_key": "mental_health",
        "keywords": [
            "mental health",
            "behavioral",
            "counselor",
            "therapist",
            "psycholog",
            "psychiatr",
            "substance abuse",
        ],
    },
    # Media & Entertainment (NAICS 71)
    "media_entertainment": {
        "naics": "71",
        "sector": "Media & Entertainment",
        "bls_sector": "Leisure and Hospitality",
        "talent_profile": "Creative & Media Workforce",
        "legacy_key": "media_entertainment",
        "keywords": [
            "media",
            "entertainment",
            "film",
            "broadcast",
            "streaming",
            "gaming",
            "music",
            "content",
        ],
    },
    # Food & Beverage (NAICS 72 sub / 31)
    "food_beverage": {
        "naics": "31",
        "sector": "Food & Beverage",
        "bls_sector": "Manufacturing",
        "talent_profile": "Food Service & Production Workforce",
        "legacy_key": "food_beverage",
        "keywords": [
            "food",
            "beverage",
            "brewery",
            "distillery",
            "catering",
            "bakery",
            "meat",
            "dairy",
        ],
    },
}

# Mapping from legacy industry keys (used in channels_db, research.py) to NAICS map keys
# Priority entries that should NOT be overridden by later NAICS map entries
_LEGACY_PRIORITY = {
    "general_entry_level": "general",
    "logistics_supply_chain": "transportation",
    "automotive": "manufacturing",
    "food_beverage": "food_beverage",
}
_LEGACY_TO_NAICS_KEY = {}
for _nk, _nv in INDUSTRY_NAICS_MAP.items():
    _lk = _nv.get("legacy_key") or ""
    if _lk and _lk not in _LEGACY_TO_NAICS_KEY:
        _LEGACY_TO_NAICS_KEY[_lk] = _nk
# Apply priority overrides
for _lk, _nk in _LEGACY_PRIORITY.items():
    if _nk in INDUSTRY_NAICS_MAP:
        _LEGACY_TO_NAICS_KEY[_lk] = _nk


def classify_industry(
    raw_industry: str, company_name: str = "", roles: list = None
) -> dict:
    """
    Classify an industry input into the correct NAICS sector with talent profile.
    Uses fuzzy keyword matching against the input string AND company name AND roles.

    Handles both:
    - Legacy keys (e.g., "healthcare_medical", "tech_engineering") from the frontend dropdown
    - Free-text industry names (e.g., "Technology", "Healthcare") from API/manual input

    Returns the full industry profile dict including legacy_key for backward compatibility.
    """
    if not raw_industry:
        raw_industry = ""

    # Step 1: Check if the input is already a legacy key (from frontend dropdown)
    raw_stripped = raw_industry.strip()
    if raw_stripped in _LEGACY_TO_NAICS_KEY:
        return INDUSTRY_NAICS_MAP[_LEGACY_TO_NAICS_KEY[raw_stripped]]

    # Step 2: Check if the input directly matches a NAICS map key
    raw_lower = raw_stripped.lower()
    if raw_lower in INDUSTRY_NAICS_MAP:
        return INDUSTRY_NAICS_MAP[raw_lower]

    # Step 3: Role-title-based industry inference (when no explicit industry)
    # This prevents defaulting to "Retail" / generic when roles clearly indicate
    # a specific industry (e.g., "Software Engineer" -> Technology).
    _ROLE_INDUSTRY_MAP: dict[str, str] = {
        # Technology
        "software engineer": "technology",
        "developer": "technology",
        "devops": "technology",
        "sre": "technology",
        "site reliability": "technology",
        "frontend": "technology",
        "backend": "technology",
        "full stack": "technology",
        "fullstack": "technology",
        "platform engineer": "technology",
        "cloud engineer": "technology",
        "security engineer": "technology",
        "qa engineer": "technology",
        "test engineer": "technology",
        "mobile developer": "technology",
        "ios developer": "technology",
        "android developer": "technology",
        "solutions architect": "technology",
        # Technology / Data
        "data scientist": "technology",
        "ml engineer": "technology",
        "machine learning": "technology",
        "data engineer": "technology",
        "data analyst": "technology",
        "analytics engineer": "technology",
        "ai engineer": "technology",
        "artificial intelligence": "technology",
        # Healthcare
        "nurse": "healthcare",
        "registered nurse": "healthcare",
        "doctor": "healthcare",
        "physician": "healthcare",
        "pharmacist": "healthcare",
        "medical": "healthcare",
        "therapist": "healthcare",
        "clinical": "healthcare",
        "radiologist": "healthcare",
        "surgeon": "healthcare",
        "dental": "healthcare",
        "paramedic": "healthcare",
        "emt": "healthcare",
        "phlebotomist": "healthcare",
        "medical assistant": "healthcare",
        # Finance
        "accountant": "finance",
        "financial analyst": "finance",
        "controller": "finance",
        "auditor": "finance",
        "actuary": "finance",
        "underwriter": "finance",
        "loan officer": "finance",
        "investment analyst": "finance",
        "portfolio manager": "finance",
        "compliance officer": "finance",
        "financial advisor": "finance",
        # Marketing (maps to general but with marketing talent profile)
        "marketing manager": "technology",
        "content strategist": "technology",
        "seo": "technology",
        "growth": "technology",
        "digital marketing": "technology",
        # Sales (maps to general but role-aware)
        "sales": "retail",
        "account executive": "retail",
        "bdr": "retail",
        "sdr": "retail",
        "business development": "retail",
    }

    if roles:
        _role_industry_votes: dict[str, int] = {}
        for _r in roles:
            _role_str = str(_r).lower().strip() if _r else ""
            for _role_kw, _ind_key in _ROLE_INDUSTRY_MAP.items():
                if _role_kw in _role_str:
                    _role_industry_votes[_ind_key] = _role_industry_votes.get(
                        _ind_key, 0
                    ) + len(_role_kw)
                    break  # First match per role wins

        if _role_industry_votes:
            _best_role_ind = max(_role_industry_votes, key=_role_industry_votes.get)
            if _best_role_ind in INDUSTRY_NAICS_MAP:
                # Only use role-based inference when raw_industry is empty/generic
                if not raw_stripped or raw_lower in (
                    "",
                    "general",
                    "other",
                    "n/a",
                    "na",
                    "none",
                ):
                    logger.info(
                        "Industry inferred from role titles: %s (votes=%s)",
                        _best_role_ind,
                        _role_industry_votes,
                    )
                    return INDUSTRY_NAICS_MAP[_best_role_ind]

    # Step 4: Fuzzy keyword matching against input + company name + roles
    search_text = f"{raw_industry} {company_name} {' '.join(roles or [])}".lower()

    best_match = None
    best_score = 0

    for key, profile in INDUSTRY_NAICS_MAP.items():
        score = 0
        for kw in profile["keywords"]:
            if kw in search_text:
                # Longer keyword matches are weighted higher
                score += len(kw)
        if score > best_score:
            best_score = score
            best_match = profile

    if best_match and best_score >= 3:
        return best_match

    # Step 5: Fallback - try to match the raw industry string directly against sector names
    if raw_lower:  # Only attempt if we have a non-empty industry string
        for key, profile in INDUSTRY_NAICS_MAP.items():
            if key in raw_lower or raw_lower in profile["sector"].lower():
                return profile

    # Step 6: Last resort - try role-based inference even when raw_industry is set
    # (covers cases where user typed a non-matching industry string)
    if roles:
        _role_industry_votes_2: dict[str, int] = {}
        for _r in roles:
            _role_str = str(_r).lower().strip() if _r else ""
            for _role_kw, _ind_key in _ROLE_INDUSTRY_MAP.items():
                if _role_kw in _role_str:
                    _role_industry_votes_2[_ind_key] = _role_industry_votes_2.get(
                        _ind_key, 0
                    ) + len(_role_kw)
                    break
        if _role_industry_votes_2:
            _best_role_ind_2 = max(
                _role_industry_votes_2, key=_role_industry_votes_2.get
            )
            if _best_role_ind_2 in INDUSTRY_NAICS_MAP:
                logger.info(
                    "Industry inferred from role titles (last resort): %s",
                    _best_role_ind_2,
                )
                return INDUSTRY_NAICS_MAP[_best_role_ind_2]

    # Final fallback
    return {
        "naics": "00",
        "sector": (
            raw_industry.strip()
            if raw_industry and raw_industry.strip()
            else "General / Multi-Industry"
        ),
        "bls_sector": "Total Nonfarm",
        "talent_profile": "General / Mixed Workforce",
        "legacy_key": "general_entry_level",
        "keywords": [],
    }


BASE_DIR = os.path.dirname(os.path.abspath(__file__))
DATA_DIR = os.path.join(BASE_DIR, "data")

# Server start time for uptime tracking (used by /health endpoint)
_SERVER_START_TIME = time.time()

# ===============================================================================
# BLUE-GREEN DEPLOY TRACKING (Linear JOV-21)
# ===============================================================================
# Zero-downtime deployment strategy for Render.com:
#
# Architecture:
#   Render.com Standard tier deploys by spinning up a NEW instance, running
#   the health check, and only routing traffic once it passes. This is
#   effectively blue-green: the old instance (blue) serves traffic while the
#   new instance (green) warms up. Render tears down blue after green is healthy.
#
# Warmup sequence:
#   1. Server starts, binds port, begins accepting connections
#   2. _deferred_startup() thread loads KB, builds vector index, inits modules
#   3. Once complete, _DEPLOY_WARMUP_COMPLETE is set to True
#   4. /api/deploy/ready returns 200 ONLY when warmup is complete
#   5. Render healthCheckPath (/) returns 200 immediately (liveness)
#   6. /api/deploy/ready is the READINESS gate -- use in CI/CD scripts
#
# Deploy version tracking:
#   RENDER_GIT_COMMIT is set automatically by Render on each deploy.
#   _DEPLOY_VERSION combines the app version with the commit hash for
#   traceability across blue/green instances.
#
# Usage in CI/CD:
#   After push, poll: curl -sf https://YOUR_APP/api/deploy/ready
#   Returns 200 with {"ready": true} when the new instance is fully warmed.
#   Returns 503 with {"ready": false} while still warming up.
# ===============================================================================

_DEPLOY_WARMUP_COMPLETE: bool = False
_DEPLOY_GIT_COMMIT: str = os.environ.get("RENDER_GIT_COMMIT") or "local"
_DEPLOY_VERSION: str = f"4.0.0-{_DEPLOY_GIT_COMMIT[:8]}"
_DEPLOY_INSTANCE_ID: str = os.environ.get("RENDER_INSTANCE_ID") or "local"

TEMPLATES_DIR = os.path.join(BASE_DIR, "templates")

# --- Request logging helpers ---
REQUEST_LOG_FILE = os.path.join(DATA_DIR, "request_log.json")
REQUEST_LOG_LOCK = REQUEST_LOG_FILE + ".lock"


def load_request_log():
    try:
        with open(REQUEST_LOG_FILE, "r") as f:
            return json.load(f)
    except (FileNotFoundError, json.JSONDecodeError):
        return []


def save_request_log(log):
    with open(REQUEST_LOG_FILE, "w") as f:
        json.dump(log, f, indent=2, default=str)


def _cleanup_old_docs(docs_dir, max_files=200):
    """Remove oldest generated docs if directory exceeds max_files."""
    try:
        files = sorted(
            [f for f in os.listdir(docs_dir) if f.endswith(".zip")],
            key=lambda f: os.path.getmtime(os.path.join(docs_dir, f)),
        )
        while len(files) > max_files:
            oldest = files.pop(0)
            os.remove(os.path.join(docs_dir, oldest))
    except OSError:
        pass


def log_request(
    data, status, file_size=0, generation_time=0, error_msg=None, doc_filename=None
):
    """Append a log entry with file locking and persist to Supabase plan_events table.

    Writes to both the local JSON file (backward compat) and the Supabase
    ``plan_events`` table for durable storage that survives deploys.
    """
    # Determine event_type from status
    event_type_map = {
        "success": "plan_generated",
        "error": "plan_failed",
        "started": "plan_started",
    }
    event_type = event_type_map.get(status, f"plan_{status}")

    # Extract enrichment APIs list
    enrichment_apis = (
        data.get("_enriched", {}).get("enrichment_summary", {}).get("apis_succeeded")
        or []
        if isinstance(data.get("_enriched"), dict)
        else []
    )

    # Extract channels selected from data
    channels_selected = data.get("channels") or data.get("selected_channels") or []

    # Build local log entry (backward compat)
    entry = {
        "id": str(uuid.uuid4())[:8],
        "timestamp": datetime.datetime.now().isoformat(),
        "requester_name": data.get("requester_name", "Unknown"),
        "requester_email": data.get("requester_email", "Unknown"),
        "client_name": data.get("client_name") or data.get("company_name", "Unknown"),
        "industry": data.get("industry", "Unknown"),
        "budget": data.get("budget", "Not specified"),
        "roles": data.get("target_roles") or data.get("roles") or [],
        "locations": data.get("locations") or [],
        "work_environment": data.get("work_environment", "Not specified"),
        "status": status,
        "file_size_bytes": file_size,
        "generation_time_seconds": round(generation_time, 2),
        "error_message": error_msg,
        "doc_filename": doc_filename,
        "enrichment_apis": enrichment_apis,
    }

    # ── Local JSON file (backward compat) ──
    try:
        os.makedirs(os.path.dirname(REQUEST_LOG_LOCK), exist_ok=True)
        with open(REQUEST_LOG_LOCK, "w") as lock_fd:
            try:
                if fcntl:
                    fcntl.flock(lock_fd, fcntl.LOCK_EX)
                log = load_request_log()
                log.append(entry)
                # Keep only last 1000 entries to prevent unbounded growth
                if len(log) > 1000:
                    log = log[-1000:]
                save_request_log(log)
            finally:
                if fcntl:
                    fcntl.flock(lock_fd, fcntl.LOCK_UN)
    except OSError as e:
        logger.error("Failed to write request log: %s", e)

    # ── Supabase plan_events table (best-effort, never raises) ──
    try:
        from supabase_client import get_client

        client = get_client()
        if client is not None:
            supabase_row = {
                "event_type": event_type,
                "plan_id": data.get("_plan_id")
                or data.get("_request_id")
                or entry["id"],
                "client_name": entry["client_name"],
                "industry": entry["industry"],
                "budget": str(entry["budget"]),
                "roles": entry["roles"],
                "locations": entry["locations"],
                "channels_selected": channels_selected,
                "user_email": entry["requester_email"],
                "user_name": entry["requester_name"],
                "generation_time_ms": round(generation_time * 1000, 1),
                "file_size_bytes": file_size,
                "error_message": error_msg,
                "enrichment_apis": enrichment_apis,
                "metadata": {
                    "work_environment": data.get("work_environment", ""),
                    "doc_filename": doc_filename,
                    "deploy_version": _DEPLOY_VERSION,
                },
            }
            client.table("plan_events").insert(supabase_row).execute()
            logger.debug("Plan event persisted to Supabase: %s", event_type)
    except ImportError:
        logger.debug("supabase_client not available; skipping plan_events persistence")
    except Exception as exc:
        logger.warning("Failed to persist plan event to Supabase: %s", exc)

    return entry


# Import research module for real data
sys.path.insert(0, BASE_DIR)
import research

# ── Canonical taxonomy standardizer ──
# Normalizes industry, role, location, platform, and metric names across
# the entire pipeline so that every subsystem uses the same canonical keys.
# Without this, industry/role/location strings from the frontend can arrive
# in dozens of variant forms and silently miss KB lookups downstream.
try:
    from standardizer import (
        normalize_industry as std_normalize_industry,
        normalize_role as std_normalize_role,
        normalize_location as std_normalize_location,
        CANONICAL_INDUSTRIES,
        get_soc_code as std_get_soc_code,
        get_role_tier as std_get_role_tier,
    )

    _STANDARDIZER_AVAILABLE = True
    logger.info("standardizer loaded successfully")
except ImportError as e:
    _STANDARDIZER_AVAILABLE = False
    logger.warning("standardizer import failed: %s", e)

try:
    from ppt_generator import generate_pptx, INDUSTRY_ALLOC_PROFILES

    logger.info("ppt_generator loaded successfully")
except ImportError as e:
    logger.warning("ppt_generator import failed: %s", e)
    generate_pptx = None
    INDUSTRY_ALLOC_PROFILES = None

try:
    from deck_generator import DeckGenerator, DeckGenerationError

    _deck_generator = DeckGenerator()
    logger.info("deck_generator (7-tier fallback) loaded successfully")
except ImportError as e:
    logger.warning("deck_generator import failed: %s", e)
    _deck_generator = None
    DeckGenerationError = None

try:
    from api_enrichment import enrich_data

    logger.info("api_enrichment loaded successfully")
except ImportError as e:
    logger.warning("api_enrichment import failed: %s", e)
    enrich_data = None

try:
    from data_synthesizer import synthesize as data_synthesize

    logger.info("data_synthesizer loaded successfully")
except ImportError as e:
    logger.warning("data_synthesizer import failed: %s", e)
    data_synthesize = None

try:
    from budget_engine import calculate_budget_allocation

    logger.info("budget_engine loaded successfully")
except ImportError as e:
    logger.warning("budget_engine import failed: %s", e)
    calculate_budget_allocation = None

# v3: Trend engine and collar intelligence for new Excel worksheets
try:

    _HAS_TREND_ENGINE = True
    logger.info("trend_engine loaded for app.py Excel worksheets")
except ImportError:
    _HAS_TREND_ENGINE = False

try:

    _HAS_COLLAR_INTEL = True
    logger.info("collar_intelligence loaded for app.py Excel worksheets")
except ImportError:
    _HAS_COLLAR_INTEL = False

# v3.5: Consolidated 4-sheet Excel generator (replaces 26+ sheet legacy generator)
try:
    from excel_v2 import generate_excel_v2

    logger.info("excel_v2 loaded successfully -- 4-sheet consolidated generator")
except ImportError as e:
    logger.warning("excel_v2 import failed, falling back to legacy generator: %s", e)
    generate_excel_v2 = None

# ── Google Cloud integration modules (lazy, non-blocking) ──
try:
    from google_maps_integration import (
        geocode_address,
        batch_geocode,
        places_autocomplete,
        get_status as maps_status,
    )

    _has_google_maps = True
    logger.info("google_maps_integration loaded successfully")
except ImportError:
    _has_google_maps = False
    geocode_address = batch_geocode = places_autocomplete = maps_status = None

try:
    from google_knowledge_graph import (
        enrich_company,
        enrich_job_title,
        enrich_location,
        get_status as kg_status,
    )

    _has_knowledge_graph = True
    logger.info("google_knowledge_graph loaded successfully")
except ImportError:
    _has_knowledge_graph = False
    enrich_company = enrich_job_title = enrich_location = kg_status = None

try:
    from google_bigquery_integration import (
        store_media_plan as bq_store_plan,
        ensure_tables_exist as bq_ensure_tables,
        get_status as bq_status,
    )

    _has_bigquery = True
    logger.info("google_bigquery_integration loaded successfully")
except ImportError:
    _has_bigquery = False

try:
    from google_search_pagespeed import (
        audit_career_page,
        get_status as pagespeed_status,
    )

    _has_pagespeed = True
    logger.info("google_search_pagespeed loaded successfully")
except ImportError:
    _has_pagespeed = False

try:
    from google_vision_integration import (
        extract_text_from_image,
        extract_text_from_pdf_vision,
        get_status as vision_status,
    )

    _has_google_vision = True
    logger.info("google_vision_integration loaded successfully")
except ImportError:
    _has_google_vision = False

try:
    from google_cloud_storage import upload_file as gcs_upload, get_status as gcs_status

    _has_gcs = True
    logger.info("google_cloud_storage loaded successfully")
except ImportError:
    _has_gcs = False

try:
    from google_ads_analytics import (
        get_google_ads_benchmarks,
        get_status as ads_analytics_status,
    )

    _has_ads_analytics = True
    logger.info("google_ads_analytics loaded successfully")
except ImportError:
    _has_ads_analytics = False

# ═══════════════════════════════════════════════════════════════════════════════
# ROLE-TIER CLASSIFICATION ENGINE
# ═══════════════════════════════════════════════════════════════════════════════

ROLE_TIER_KEYWORDS = {
    "executive": {
        "tier": "Executive / Leadership",
        "keywords": [
            "ceo",
            "cfo",
            "cto",
            "cio",
            "cmo",
            "vp ",
            "vice president",
            "svp",
            "evp",
            "chief",
            "president",
            "director",
            "head of",
        ],
        "funnel": {
            "awareness": 0.10,
            "consideration": 0.25,
            "engagement": 0.35,
            "conversion": 0.30,
        },
        "channels": [
            "Executive search firms",
            "LinkedIn Recruiter",
            "Board networks",
            "Industry conferences",
        ],
        "cpa_multiplier": 3.0,
        "time_to_fill_days": 90,
    },
    "professional": {
        "tier": "Professional / White-Collar",
        "keywords": [
            "engineer",
            "developer",
            "analyst",
            "manager",
            "consultant",
            "architect",
            "scientist",
            "specialist",
            "accountant",
            "auditor",
            "attorney",
            "lawyer",
            "planner",
            "strategist",
            "designer",
            "product",
            "data",
            "researcher",
            "coordinator",
            "administrator",
            "compliance",
            "underwriter",
        ],
        "funnel": {
            "awareness": 0.20,
            "consideration": 0.30,
            "engagement": 0.30,
            "conversion": 0.20,
        },
        "channels": [
            "LinkedIn",
            "Indeed",
            "Glassdoor",
            "Company careers page",
            "Industry job boards",
        ],
        "cpa_multiplier": 1.5,
        "time_to_fill_days": 45,
    },
    "clinical": {
        "tier": "Clinical / Licensed",
        "keywords": [
            "nurse",
            "rn ",
            "lpn",
            "cna",
            "physician",
            "doctor",
            "therapist",
            "pharmacist",
            "paramedic",
            "emt",
            "medical assistant",
            "dental",
            "radiology",
            "sonograph",
            "surgical tech",
            "respiratory",
        ],
        "funnel": {
            "awareness": 0.15,
            "consideration": 0.25,
            "engagement": 0.35,
            "conversion": 0.25,
        },
        "channels": [
            "NurseFly",
            "Vivian Health",
            "Health eCareers",
            "Indeed Healthcare",
            "Hospital networks",
        ],
        "cpa_multiplier": 2.0,
        "time_to_fill_days": 60,
    },
    "skilled_trades": {
        "tier": "Skilled Trades / Technical",
        "keywords": [
            "technician",
            "mechanic",
            "electrician",
            "plumber",
            "welder",
            "hvac",
            "machinist",
            "installer",
            "maintenance",
            "repair",
            "inspector",
            "quality",
            "foreman",
            "carpenter",
            "aircraft mechanic",
            "solar panel",
            "wind turbine",
        ],
        "funnel": {
            "awareness": 0.25,
            "consideration": 0.25,
            "engagement": 0.25,
            "conversion": 0.25,
        },
        "channels": [
            "Trade schools",
            "Union halls",
            "Indeed Blue",
            "Craigslist",
            "Local job fairs",
        ],
        "cpa_multiplier": 1.2,
        "time_to_fill_days": 30,
    },
    "hourly": {
        "tier": "Hourly / Entry-Level",
        "keywords": [
            "associate",
            "clerk",
            "cashier",
            "barista",
            "crew",
            "team member",
            "warehouse",
            "package handler",
            "picker",
            "packer",
            "stocker",
            "janitor",
            "custodian",
            "housekeeper",
            "dishwasher",
            "food prep",
            "line cook",
            "server",
            "host",
            "busser",
            "retail associate",
            "store associate",
            "front desk",
            "receptionist",
            "call center",
            "customer service rep",
            "csr",
        ],
        "funnel": {
            "awareness": 0.35,
            "consideration": 0.25,
            "engagement": 0.20,
            "conversion": 0.20,
        },
        "channels": [
            "Indeed",
            "Snagajob",
            "Jobcase",
            "Facebook Jobs",
            "Walk-in/Referral",
            "Text-to-Apply",
        ],
        "cpa_multiplier": 0.5,
        "time_to_fill_days": 14,
    },
    "gig": {
        "tier": "Gig / Independent Contractor",
        "keywords": [
            "driver partner",
            "delivery partner",
            "delivery driver",
            "courier",
            "rider",
            "dasher",
            "shopper",
            "freelance",
            "contractor",
            "1099",
            "independent",
            "flex",
            "on-demand",
        ],
        "funnel": {
            "awareness": 0.40,
            "consideration": 0.20,
            "engagement": 0.15,
            "conversion": 0.25,
        },
        "channels": [
            "Programmatic ads",
            "Social media (Facebook/Instagram)",
            "Referral bonuses",
            "Google Ads",
            "SMS campaigns",
        ],
        "cpa_multiplier": 0.3,
        "time_to_fill_days": 7,
    },
    "education": {
        "tier": "Education / Academic",
        "keywords": [
            "teacher",
            "professor",
            "instructor",
            "principal",
            "superintendent",
            "dean",
            "librarian",
            "counselor",
            "aide",
            "paraprofessional",
            "substitute",
            "tutor",
            "curriculum",
            "coach",
        ],
        "funnel": {
            "awareness": 0.20,
            "consideration": 0.30,
            "engagement": 0.30,
            "conversion": 0.20,
        },
        "channels": [
            "SchoolSpring",
            "K12JobSpot",
            "HigherEdJobs",
            "State education boards",
            "University career offices",
        ],
        "cpa_multiplier": 1.0,
        "time_to_fill_days": 45,
    },
}


def classify_role_tier(role_title: str) -> dict:
    """Classify a role into a tier based on keyword matching. Returns the tier profile."""
    role_lower = role_title.lower().strip()

    best_match = None
    best_score = 0

    for tier_key, tier_data in ROLE_TIER_KEYWORDS.items():
        score = 0
        for kw in tier_data["keywords"]:
            if kw in role_lower:
                score += len(kw)  # Longer keyword matches score higher
        if score > best_score:
            best_score = score
            best_match = tier_data

    if best_match and best_score >= 2:
        return best_match

    # Default to professional if no match
    return ROLE_TIER_KEYWORDS["professional"]


# Load global supply data (original 107KB file)
GLOBAL_SUPPLY_PATH = os.path.join(
    os.path.dirname(os.path.abspath(__file__)), "data", "global_supply.json"
)
global_supply_data = {}
try:
    with open(GLOBAL_SUPPLY_PATH, "r") as f:
        global_supply_data = json.load(f)
except (FileNotFoundError, json.JSONDecodeError, OSError):
    pass

# Load NEW expanded supply repository (2.7MB, added S30)
_DATA_DIR = os.path.dirname(os.path.abspath(__file__))
_expanded_supply_repo: dict = {}
try:
    _esp = os.path.join(_DATA_DIR, "data", "joveo_global_supply_repository.json")
    if os.path.exists(_esp):
        with open(_esp, "r") as f:
            _expanded_supply_repo = json.load(f)
        logger.info(
            "Loaded expanded supply repository: %d entries",
            (
                len(_expanded_supply_repo)
                if isinstance(_expanded_supply_repo, (list, dict))
                else 0
            ),
        )
except (FileNotFoundError, json.JSONDecodeError, OSError) as e:
    logger.warning(f"Could not load expanded supply repository: {e}")

# Load client plans from data/client_plans/ directory (RTX + others, added S30)
_client_plans_dir = os.path.join(_DATA_DIR, "data", "client_plans")
_client_plans_data: dict = {}
try:
    if os.path.isdir(_client_plans_dir):
        for fname in os.listdir(_client_plans_dir):
            if fname.endswith(".json"):
                fpath = os.path.join(_client_plans_dir, fname)
                with open(fpath, "r") as f:
                    plan_key = fname.replace(".json", "")
                    _client_plans_data[plan_key] = json.load(f)
        if _client_plans_data:
            logger.info(
                "Loaded %d client plan files: %s",
                len(_client_plans_data),
                ", ".join(_client_plans_data.keys()),
            )
except (FileNotFoundError, json.JSONDecodeError, OSError) as e:
    logger.warning(f"Could not load client plans directory: {e}")

_channels_db_cache = None
_joveo_publishers_cache = None


def load_channels_db():
    global _channels_db_cache
    if _channels_db_cache is not None:
        return _channels_db_cache
    try:
        with open(os.path.join(DATA_DIR, "channels_db.json"), "r") as f:
            _channels_db_cache = json.load(f)
            return _channels_db_cache
    except (FileNotFoundError, json.JSONDecodeError, OSError) as e:
        logger.error("Failed to load channels_db.json: %s", e)
        return {}


def load_joveo_publishers():
    global _joveo_publishers_cache
    if _joveo_publishers_cache is not None:
        return _joveo_publishers_cache
    path = os.path.join(DATA_DIR, "joveo_publishers.json")
    try:
        with open(path, "r", encoding="utf-8") as f:
            _joveo_publishers_cache = json.load(f)
            return _joveo_publishers_cache
    except FileNotFoundError:
        logger.warning("joveo_publishers.json not found at %s", path)
        return {}
    except json.JSONDecodeError as e:
        logger.error("joveo_publishers.json JSON parse error: %s", e)
        return {}
    except OSError as e:
        logger.error("Failed to read joveo_publishers.json: %s", e)
        return {}


def fetch_client_logo(client_name, client_website=""):
    """Try to fetch a client logo using the client website URL (preferred) or name-based domain guess."""
    domain = ""

    # 1. Try to extract domain from the provided client website URL (most accurate)
    if client_website:
        try:
            parsed = urlparse(
                client_website
                if "://" in client_website
                else f"https://{client_website}"
            )
            # SSRF prevention: only allow http/https schemes
            if parsed.scheme.lower() not in ("http", "https"):
                logger.warning(f"Blocked unsafe URL scheme: {parsed.scheme}")
                domain = ""
            else:
                domain = parsed.hostname or ""
                # Remove www. prefix for cleaner domain
                if domain.startswith("www."):
                    domain = domain[4:]
        except ValueError:
            pass

    # Block private/internal IPs (SSRF prevention)
    if domain:
        _blocked_patterns = (
            "localhost",
            "127.0.0.1",
            "0.0.0.0",
            "169.254.",
            "10.",
            "172.16.",
            "172.17.",
            "172.18.",
            "172.19.",
            "172.20.",
            "172.21.",
            "172.22.",
            "172.23.",
            "172.24.",
            "172.25.",
            "172.26.",
            "172.27.",
            "172.28.",
            "172.29.",
            "172.30.",
            "172.31.",
            "192.168.",
            "[::1]",
            "metadata.google",
        )
        if any(
            domain.startswith(p) or domain == p.rstrip(".") for p in _blocked_patterns
        ):
            logger.warning(f"Blocked internal/private domain: {domain}")
            domain = ""

    # 2. Fallback: guess domain from client name
    if not domain and client_name:
        name = client_name.lower().strip()
        for suffix in [
            " inc",
            " inc.",
            " llc",
            " ltd",
            " corp",
            " corporation",
            " co",
            " company",
            " group",
            " international",
        ]:
            name = name.replace(suffix, "")
        domain = re.sub(r"[^a-z0-9]", "", name) + ".com"

    if not domain:
        return None, None

    # Try Clearbit Logo API (high quality, free for logos), then Google favicon
    logo_urls = [
        f"https://logo.clearbit.com/{domain}",
        f"https://www.google.com/s2/favicons?domain={domain}&sz=128",
    ]

    for url in logo_urls:
        try:
            req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
            response = urllib.request.urlopen(req, timeout=5)
            if response.status == 200:
                img_data = response.read()
                if len(img_data) > 500:  # Ensure it's not a tiny placeholder
                    return img_data, url
        except Exception:
            continue
    return None, None


# ═══════════════════════════════════════════════════════════════════════════════
# INDUSTRY-SPECIFIC NICHE CHANNELS MAPPING
# ═══════════════════════════════════════════════════════════════════════════════
INDUSTRY_NICHE_CHANNELS = {
    "technology": [
        "Stack Overflow Jobs",
        "GitHub Jobs",
        "HackerNews (Who's Hiring)",
        "AngelList/Wellfound",
        "Dice",
        "Built In",
        "Triplebyte",
    ],
    "tech_engineering": [
        "Stack Overflow Jobs",
        "GitHub Jobs",
        "HackerNews (Who's Hiring)",
        "AngelList/Wellfound",
        "Dice",
        "Built In",
        "Triplebyte",
    ],
    "healthcare": [
        "NurseFly",
        "Vivian Health",
        "Health eCareers",
        "Doximity",
        "PracticeLink",
        "Medscape Career Center",
        "AllNurses",
    ],
    "healthcare_medical": [
        "NurseFly",
        "Vivian Health",
        "Health eCareers",
        "Doximity",
        "PracticeLink",
        "Medscape Career Center",
        "AllNurses",
    ],
    "finance": [
        "eFinancialCareers",
        "Wall Street Oasis",
        "Selby Jennings",
        "The Muse (Finance)",
        "Vault.com",
    ],
    "finance_banking": [
        "eFinancialCareers",
        "Wall Street Oasis",
        "Selby Jennings",
        "The Muse (Finance)",
        "Vault.com",
    ],
    "hospitality": ["Hcareers", "Poached", "Culinary Agents", "Harri", "OysterLink"],
    "hospitality_travel": [
        "Hcareers",
        "Poached",
        "Culinary Agents",
        "Harri",
        "OysterLink",
    ],
    "education": [
        "SchoolSpring",
        "K12JobSpot",
        "HigherEdJobs",
        "Chronicle Vitae",
        "TeachAway",
        "Teach.org",
    ],
    "manufacturing": [
        "ManufacturingJobs.com",
        "iHireManufacturing",
        "JobsInManufacturing",
        "IndustryWeek Careers",
    ],
    "transportation": [
        "CDLJobs.com",
        "TruckingJobs.com",
        "FleetOwner Careers",
        "Transport Topics Jobs",
        "DriveForMe",
    ],
    "energy": [
        "Rigzone",
        "EnergyJobline",
        "Solar Energy Jobs",
        "NABCEP Job Board",
        "Oil and Gas Job Search",
    ],
    "energy_utilities": [
        "Rigzone",
        "EnergyJobline",
        "Solar Energy Jobs",
        "NABCEP Job Board",
        "Oil and Gas Job Search",
    ],
    "retail": ["RetailCareersNow", "AllRetailJobs", "iHireRetail", "NRF Job Board"],
    "retail_consumer": [
        "RetailCareersNow",
        "AllRetailJobs",
        "iHireRetail",
        "NRF Job Board",
    ],
    "professional_services": [
        "Vault.com",
        "Management Consulted",
        "Fishbowl",
        "WetFeet",
        "Idealist (for nonprofits)",
    ],
    "construction": [
        "ConstructionJobs.com",
        "iHireConstruction",
        "AGC Career Center",
        "BuildingTeams",
    ],
    "construction_real_estate": [
        "ConstructionJobs.com",
        "iHireConstruction",
        "AGC Career Center",
        "BuildingTeams",
    ],
    "rideshare": [
        "Ridester",
        "The Rideshare Guy",
        "Gridwise community",
        "Driver forums",
        "Facebook gig groups",
    ],
    "government": [
        "USAJobs",
        "GovernmentJobs.com",
        "ClearanceJobs",
        "Intelligence Careers",
    ],
    "nonprofit": [
        "Idealist",
        "Work for Good",
        "DevEx",
        "NGOJobBoard",
        "Foundation List",
    ],
    "blue_collar_trades": [
        "TradesmenInternational",
        "CraftHire",
        "PeopleReady",
        "iHireConstruction",
        "CDLJobs.com",
    ],
    "maritime_marine": [
        "RigZone",
        "SeaCareer",
        "MarineLink Jobs",
        "AllMarinerJobs",
        "gCaptain Jobs",
    ],
    "military_recruitment": [
        "RecruitMilitary",
        "Military.com",
        "ClearanceJobs",
        "Hire Heroes USA",
        "USAJobs.gov",
    ],
    "legal_services": [
        "LawJobs.com",
        "Robert Half Legal",
        "Above the Law",
        "NALP Job Board",
        "Lawcrossing.com",
    ],
    "mental_health": [
        "Psychology Today Jobs",
        "NASW Career Center",
        "APA PsycCareers",
        "SocialWorkJobBank",
        "iHireTherapy",
    ],
    "aerospace_defense": [
        "ClearanceJobs",
        "ClearedJobs.Net",
        "Aviation Job Search",
        "SpaceCareers.uk",
        "Airswift",
    ],
    "pharma_biotech": [
        "PharmiWeb.com",
        "BioSpace",
        "MedReps",
        "Science Careers (AAAS)",
        "Nature Careers",
    ],
    "insurance": [
        "InsuranceJobs.com",
        "GreatInsuranceJobs.com",
        "Actuarial Careers Inc",
        "SOA Career Center",
    ],
    "telecommunications": [
        "WirelessEstimator Jobs",
        "FierceTelecom Careers",
        "Light Reading Jobs",
        "IEEE Job Site",
    ],
    "automotive": [
        "AutomotiveJobFinder",
        "iHireEngineering",
        "SAE International Jobs",
        "AutoCareers",
    ],
    "food_beverage": [
        "FoodIndustryJobs.com",
        "iHireHospitality",
        "FoodProcessing.com Careers",
        "IFT Career Center",
    ],
    "logistics_supply_chain": [
        "SupplyChainRecruit.com",
        "LogisticsJobShop.com",
        "3PL Jobs",
        "WarehouseJobs.com",
    ],
    "media_entertainment": [
        "ProductionHub",
        "Mandy.com",
        "EntertainmentCareers.net",
        "MediaBistro",
    ],
    "general_entry_level": [
        "Snagajob",
        "Wonolo",
        "Instawork",
        "Jobcase",
        "College Recruiter",
    ],
}


def _verify_plan_data(data):
    """Use Gemini/secondary LLM to verify key facts in plan data before output.

    Checks channel recommendations, benchmark numbers, and location appropriateness.
    Returns verification status dict. Non-blocking -- failures return 'skipped'.
    """
    try:
        from llm_router import call_llm, TASK_VERIFICATION
    except ImportError:
        return {"status": "skipped", "reason": "llm_router_unavailable"}

    synthesized = data.get("_synthesized", {})
    if not isinstance(synthesized, dict):
        return {"status": "skipped", "reason": "no_synthesized_data"}

    # Build a concise snapshot of key plan data points for verification
    client = data.get("client_name", "Client")
    industry = data.get("industry_label", data.get("industry") or "")
    locations = data.get("locations") or []
    roles = data.get("roles") or []
    budget = data.get("budget", "N/A")
    budget_alloc = data.get("_budget_allocation", {})

    # Extract key numbers to verify
    check_items = []
    if isinstance(budget_alloc, dict):
        meta = budget_alloc.get("metadata", {})
        if meta:
            check_items.append(f"Total budget: ${meta.get('total_budget', 'N/A')}")
        ch_allocs = budget_alloc.get("channel_allocations", {})
        for ch_name, ch_data in list(ch_allocs.items())[:5]:
            if isinstance(ch_data, dict):
                check_items.append(
                    f"{ch_name}: ${ch_data.get('dollar_amount') or 0:.0f}, "
                    f"CPC=${ch_data.get('cpc') or 0:.2f}, "
                    f"est clicks={ch_data.get('projected_clicks') or 0}"
                )

    if not check_items:
        return {"status": "skipped", "reason": "no_data_to_verify"}

    # ── Cross-reference with KB benchmark data (Priority 3) ──
    # Pull Appcast 2026 and Google Ads 2025 benchmarks for LLM verification context
    kb_benchmark_context = ""
    try:
        _vkb = load_knowledge_base()
        _vwp = _vkb.get("white_papers", {})
        _vappcast = (
            _vwp.get("reports", {})
            .get("appcast_benchmark_2026", {})
            .get("benchmarks", {})
        )
        _vgads = _vkb.get("google_ads_benchmarks", {})

        # Appcast occupation-level benchmarks for this industry
        _V_APPCAST_MAP = {
            "healthcare_medical": "healthcare",
            "tech_engineering": "technology",
            "technology": "technology",
            "retail_consumer": "retail",
            "finance_banking": "finance",
            "logistics_supply_chain": "warehousing_logistics",
            "hospitality_travel": "hospitality",
            "manufacturing": "manufacturing",
            "general_entry_level": "administration",
        }
        _v_occ = _V_APPCAST_MAP.get(data.get("industry") or "", "")
        if _vappcast and _v_occ:
            _v_cpa = _vappcast.get("cpa_by_occupation_2025", {}).get(_v_occ, "N/A")
            _v_cph = _vappcast.get("cph_by_occupation_2025", {}).get(_v_occ, "N/A")
            _v_ar = _vappcast.get("apply_rate_by_occupation_2025", {}).get(
                _v_occ, "N/A"
            )
            kb_benchmark_context += f"\nAppcast 2026 benchmarks for {_v_occ}: CPA={_v_cpa}, CPH={_v_cph}, Apply Rate={_v_ar}"
            kb_benchmark_context += f"\nOverall medians: CPC={_vappcast.get('overall_median_cpc', 'N/A')}, CPA={_vappcast.get('overall_median_cpa', 'N/A')}, CPH={_vappcast.get('overall_median_cph', 'N/A')}"

        # Google Ads first-party CPC benchmarks
        _V_GADS_MAP = {
            "healthcare_medical": "skilled_healthcare",
            "healthcare": "skilled_healthcare",
            "pharma_biotech": "skilled_healthcare",
            "tech_engineering": "software_tech",
            "technology": "software_tech",
            "logistics_supply_chain": "logistics_supply_chain",
            "logistics": "logistics_supply_chain",
            "transportation": "logistics_supply_chain",
            "manufacturing": "logistics_supply_chain",
            "retail_consumer": "retail_hospitality",
            "retail": "retail_hospitality",
            "hospitality": "retail_hospitality",
            "hospitality_travel": "retail_hospitality",
            "finance": "corporate_professional",
            "finance_banking": "corporate_professional",
            "general_entry_level": "general_recruitment",
            "general": "general_recruitment",
            "education": "education_public_service",
        }
        _v_gcat = _V_GADS_MAP.get(data.get("industry") or "", "")
        _v_gads_cat = _vgads.get("categories", {}).get(_v_gcat, {})
        if _v_gads_cat:
            kb_benchmark_context += f"\nJoveo Google Ads 2025 ({_v_gcat}): blended CPC=${_v_gads_cat.get('blended_cpc', 'N/A')}, median CPC=${_v_gads_cat.get('cpc_stats', {}).get('median', 'N/A')}, CTR={_v_gads_cat.get('blended_ctr', 'N/A')}%"
    except Exception:
        pass  # Non-fatal

    prompt = f"""Verify this recruitment media plan data for reasonableness.

Client: {client}
Industry: {industry}
Locations: {', '.join(str(l) for l in locations[:5])}
Roles: {', '.join(str(r) for r in roles[:5])}
Budget: {budget}

Key allocations:
{chr(10).join(check_items)}

KB Benchmark Reference Data (use these to cross-check):
{kb_benchmark_context if kb_benchmark_context else "No KB benchmarks available for this industry."}

Check:
1. Are CPC values reasonable for the industry/channels? Cross-reference against the KB benchmark data above.
2. Are CPA projections consistent with Appcast occupation-level benchmarks shown above?
3. Are click/application projections mathematically consistent with budget and CPC?
4. Are channel recommendations appropriate for this industry?
5. If any generated metric deviates by more than 50% from the KB benchmark, flag it.

Return ONLY valid JSON:
{{"verified": true, "issues": [], "severity": "none"}}
OR
{{"verified": false, "issues": ["issue1", "issue2"], "severity": "minor|major"}}"""

    try:
        # S50: 10s timeout for verification -- structured JSON, Gemini preferred.
        # Falls back through TASK_VERIFICATION chain if Gemini unavailable.
        result = call_llm(
            messages=[{"role": "user", "content": prompt}],
            system_prompt="You are a recruitment advertising data verifier. Return ONLY valid JSON.",
            max_tokens=512,
            task_type=TASK_VERIFICATION,
            query_text="verify media plan data",
            preferred_providers=["gemini"],
            timeout_budget=10.0,
        )
        if result and (result.get("text") or result.get("content")):
            content = result.get("text") or result.get("content") or ""
            import re as _re

            json_match = _re.search(r"\{[\s\S]*?\}", content)
            if json_match:
                parsed = json.loads(json_match.group())
                return {
                    "status": "verified" if parsed.get("verified") else "issues_found",
                    "issues": parsed.get("issues") or [],
                    "severity": parsed.get("severity", "none"),
                    "provider": result.get("provider", "unknown"),
                }
    except Exception:
        pass

    return {"status": "skipped", "reason": "verification_failed"}


# ═══════════════════════════════════════════════════════════════════════════════
# LEGACY EXCEL GENERATOR (archived to archive/excel_legacy.py)
# ═══════════════════════════════════════════════════════════════════════════════
try:
    from archive.excel_legacy import generate_excel  # noqa: E402

    logger.info("Legacy Excel generator loaded from archive/")
except ImportError:
    logger.warning(
        "Legacy excel_legacy not found in archive/ -- legacy fallback disabled"
    )
    generate_excel = None  # type: ignore[assignment]


ADMIN_API_KEY = (
    os.environ.get("ADMIN_API_KEY") or os.environ.get("NOVA_ADMIN_KEY") or ""
)
if not ADMIN_API_KEY:
    logger.warning(
        "ADMIN_API_KEY/NOVA_ADMIN_KEY not set -- admin endpoints will reject all requests"
    )

# ═══════════════════════════════════════════════════════════════════════════════
# OPENAPI 3.0 SPECIFICATION
# ═══════════════════════════════════════════════════════════════════════════════
_OPENAPI_SPEC = {
    "openapi": "3.0.3",
    "info": {
        "title": "AI Media Planner API",
        "description": "Recruitment advertising media plan generator with AI-powered research, "
        "budget optimization, and strategy deck creation.",
        "version": "4.0.0",
        "contact": {
            "name": "Joveo Engineering",
            "url": "https://media-plan-generator.onrender.com",
        },
        "license": {"name": "Proprietary"},
    },
    "servers": [
        {
            "url": "https://media-plan-generator.onrender.com",
            "description": "Production",
        },
        {"url": "http://localhost:5001", "description": "Local development"},
    ],
    "paths": {
        "/api/generate": {
            "post": {
                "summary": "Generate media plan",
                "description": "Generate a recruitment advertising media plan bundle (Excel + PPT in ZIP). "
                "Supports synchronous (default) and asynchronous (X-Async: true) modes.",
                "operationId": "generateMediaPlan",
                "tags": ["Generation"],
                "parameters": [
                    {
                        "name": "X-Async",
                        "in": "header",
                        "required": False,
                        "schema": {"type": "string", "enum": ["true"]},
                        "description": "Set to 'true' to run generation asynchronously. "
                        "Returns a job_id for polling via /api/jobs/{job_id}.",
                    }
                ],
                "requestBody": {
                    "required": True,
                    "content": {
                        "application/json": {
                            "schema": {
                                "type": "object",
                                "required": ["client_name"],
                                "properties": {
                                    "client_name": {
                                        "type": "string",
                                        "description": "Company or client name",
                                        "example": "Acme Corp",
                                    },
                                    "industry": {
                                        "type": "string",
                                        "description": "Industry vertical (e.g. healthcare, technology)",
                                        "example": "healthcare",
                                    },
                                    "budget": {
                                        "type": "string",
                                        "description": "Campaign budget (e.g. '$50,000', '100000')",
                                        "example": "$50,000",
                                    },
                                    "roles": {
                                        "type": "array",
                                        "items": {"type": "string"},
                                        "description": "Target roles to hire",
                                        "example": [
                                            "Registered Nurse",
                                            "Medical Assistant",
                                        ],
                                    },
                                    "target_roles": {
                                        "type": "array",
                                        "items": {"type": "string"},
                                        "description": "Alias for roles",
                                    },
                                    "locations": {
                                        "type": "array",
                                        "items": {"type": "string"},
                                        "description": "Target hiring locations",
                                        "example": ["Dallas, TX", "Houston, TX"],
                                    },
                                    "competitors": {
                                        "type": "array",
                                        "items": {"type": "string"},
                                        "description": "Competitor company names for benchmarking",
                                    },
                                    "campaign_duration": {
                                        "type": "string",
                                        "description": "Campaign timeline",
                                        "example": "3-6 months",
                                    },
                                    "hire_volume": {
                                        "type": "string",
                                        "description": "Number of hires expected",
                                    },
                                    "notes": {
                                        "type": "string",
                                        "description": "Additional context or requirements",
                                    },
                                },
                            }
                        }
                    },
                },
                "responses": {
                    "200": {
                        "description": "Synchronous: Binary ZIP containing Excel + PPT. Async: JSON with job_id.",
                        "content": {
                            "application/zip": {
                                "schema": {"type": "string", "format": "binary"},
                            },
                            "application/json": {
                                "schema": {
                                    "type": "object",
                                    "properties": {
                                        "job_id": {"type": "string"},
                                        "status": {
                                            "type": "string",
                                            "enum": ["processing"],
                                        },
                                        "poll_url": {"type": "string"},
                                    },
                                },
                            },
                        },
                    },
                    "400": {
                        "description": "Invalid input (missing client_name, bad JSON, etc.)"
                    },
                    "413": {"description": "Request body too large (>10MB)"},
                    "429": {"description": "Rate limit exceeded"},
                },
            }
        },
        "/api/chat": {
            "post": {
                "summary": "Nova AI chat",
                "description": "Send a message to the Nova recruitment intelligence chatbot.",
                "operationId": "novaChat",
                "tags": ["Chat"],
                "requestBody": {
                    "required": True,
                    "content": {
                        "application/json": {
                            "schema": {
                                "type": "object",
                                "required": ["message"],
                                "properties": {
                                    "message": {
                                        "type": "string",
                                        "description": "User message",
                                        "example": "What are the best job boards for healthcare?",
                                    },
                                    "conversation_history": {
                                        "type": "array",
                                        "items": {
                                            "type": "object",
                                            "properties": {
                                                "role": {
                                                    "type": "string",
                                                    "enum": ["user", "assistant"],
                                                },
                                                "content": {"type": "string"},
                                            },
                                        },
                                        "description": "Previous conversation turns for context",
                                    },
                                },
                            }
                        }
                    },
                },
                "responses": {
                    "200": {
                        "description": "Chat response",
                        "content": {
                            "application/json": {
                                "schema": {
                                    "type": "object",
                                    "properties": {
                                        "response": {"type": "string"},
                                        "sources": {
                                            "type": "array",
                                            "items": {"type": "string"},
                                        },
                                        "confidence": {
                                            "type": "number",
                                            "format": "float",
                                        },
                                        "tools_used": {
                                            "type": "array",
                                            "items": {"type": "string"},
                                        },
                                    },
                                },
                            }
                        },
                    },
                    "429": {"description": "Rate limit exceeded"},
                },
            }
        },
        "/api/chat/stream": {
            "post": {
                "summary": "Nova AI chat (SSE streaming)",
                "description": "Send a message and receive a streaming SSE response, token by token.",
                "operationId": "novaChatStream",
                "tags": ["Chat"],
                "requestBody": {
                    "required": True,
                    "content": {
                        "application/json": {
                            "schema": {
                                "type": "object",
                                "required": ["message"],
                                "properties": {
                                    "message": {"type": "string"},
                                    "conversation_id": {"type": "string"},
                                    "history": {
                                        "type": "array",
                                        "items": {"type": "object"},
                                    },
                                },
                            }
                        }
                    },
                },
                "responses": {
                    "200": {
                        "description": "SSE event stream (Content-Type: text/event-stream)"
                    },
                    "429": {"description": "Rate limit exceeded"},
                },
            }
        },
        "/ws/chat": {
            "get": {
                "summary": "WebSocket chat streaming",
                "description": "WebSocket upgrade endpoint for real-time chat streaming. "
                "Preferred over SSE for lower latency. Clients should fall back "
                "to /api/chat/stream (SSE) if WebSocket is unavailable.",
                "operationId": "novaChatWebSocket",
                "tags": ["Chat"],
                "responses": {
                    "101": {"description": "Switching Protocols (WebSocket upgrade)"},
                    "400": {"description": "Invalid WebSocket handshake"},
                },
            }
        },
        "/api/chat/stop": {
            "post": {
                "summary": "Cancel active streaming chat",
                "description": "Stop an active SSE streaming response by conversation_id.",
                "operationId": "novaChatStop",
                "tags": ["Chat"],
                "requestBody": {
                    "required": True,
                    "content": {
                        "application/json": {
                            "schema": {
                                "type": "object",
                                "required": ["conversation_id"],
                                "properties": {"conversation_id": {"type": "string"}},
                            }
                        }
                    },
                },
                "responses": {"200": {"description": "Cancellation result"}},
            }
        },
        "/api/chat/feedback": {
            "post": {
                "summary": "Submit message feedback",
                "description": "Rate a Nova chat message with thumbs up or down.",
                "operationId": "novaChatFeedback",
                "tags": ["Chat"],
                "requestBody": {
                    "required": True,
                    "content": {
                        "application/json": {
                            "schema": {
                                "type": "object",
                                "required": ["message_id", "rating"],
                                "properties": {
                                    "message_id": {"type": "string"},
                                    "rating": {
                                        "type": "string",
                                        "enum": ["up", "down"],
                                    },
                                    "conversation_id": {"type": "string"},
                                },
                            }
                        }
                    },
                },
                "responses": {"200": {"description": "Feedback recorded"}},
            }
        },
        "/api/chat/history": {
            "get": {
                "summary": "Load conversation history",
                "description": "Load all messages for a conversation by ID from Supabase.",
                "operationId": "novaChatHistory",
                "tags": ["Chat"],
                "parameters": [
                    {
                        "name": "conversation_id",
                        "in": "query",
                        "required": True,
                        "schema": {"type": "string"},
                    }
                ],
                "responses": {"200": {"description": "Conversation messages"}},
            }
        },
        "/api/chat/conversations": {
            "get": {
                "summary": "List recent conversations",
                "description": "List the most recent Nova chat conversations (last 50 by default).",
                "operationId": "novaChatConversations",
                "tags": ["Chat"],
                "parameters": [
                    {
                        "name": "limit",
                        "in": "query",
                        "required": False,
                        "schema": {"type": "integer", "default": 50},
                    }
                ],
                "responses": {"200": {"description": "List of conversations"}},
            }
        },
        "/api/health/ping": {
            "get": {
                "summary": "Ultra-lightweight liveness ping",
                "description": "Returns immediately with no deep checks. Ideal for monitoring tools, HEAD requests, and uptime services that need sub-second responses.",
                "operationId": "healthPing",
                "tags": ["Health"],
                "responses": {
                    "200": {
                        "description": "Service is alive",
                        "content": {
                            "application/json": {
                                "schema": {
                                    "type": "object",
                                    "properties": {
                                        "status": {"type": "string", "example": "ok"},
                                        "timestamp": {
                                            "type": "string",
                                            "format": "date-time",
                                        },
                                    },
                                },
                            }
                        },
                    }
                },
            }
        },
        "/api/health": {
            "get": {
                "summary": "Liveness probe",
                "description": "Lightweight health check for load balancers and uptime monitors.",
                "operationId": "healthLiveness",
                "tags": ["Health"],
                "responses": {
                    "200": {
                        "description": "Service is alive",
                        "content": {
                            "application/json": {
                                "schema": {
                                    "type": "object",
                                    "properties": {
                                        "status": {"type": "string", "example": "ok"},
                                        "version": {
                                            "type": "string",
                                            "example": "4.0.0",
                                        },
                                        "timestamp": {
                                            "type": "string",
                                            "format": "date-time",
                                        },
                                    },
                                },
                            }
                        },
                    }
                },
            }
        },
        "/api/health/ready": {
            "get": {
                "summary": "Readiness probe",
                "description": "Deep readiness check (KB loaded, disk, memory, modules).",
                "operationId": "healthReadiness",
                "tags": ["Health"],
                "responses": {
                    "200": {"description": "Service is ready"},
                    "503": {"description": "Service is not ready (degraded)"},
                },
            }
        },
        "/api/health/data-matrix": {
            "get": {
                "summary": "Data matrix status",
                "description": "Background data matrix health monitor status.",
                "operationId": "healthDataMatrix",
                "tags": ["Health"],
                "security": [{"BearerAuth": []}],
                "responses": {
                    "200": {"description": "Data matrix status"},
                    "401": {"description": "Unauthorized"},
                    "503": {"description": "Data matrix degraded or unavailable"},
                },
            }
        },
        "/api/health/auto-qc": {
            "get": {
                "summary": "Auto QC status",
                "description": "Autonomous QC engine test results and self-upgrade status.",
                "operationId": "healthAutoQC",
                "tags": ["Health"],
                "security": [{"BearerAuth": []}],
                "responses": {
                    "200": {"description": "QC engine status"},
                    "401": {"description": "Unauthorized"},
                    "503": {"description": "QC engine degraded or unavailable"},
                },
            }
        },
        "/api/health/enrichment": {
            "get": {
                "summary": "Data enrichment status",
                "description": "Background data enrichment engine status, freshness checks, and recent activity.",
                "operationId": "healthEnrichment",
                "tags": ["Health"],
                "security": [{"BearerAuth": []}],
                "responses": {
                    "200": {"description": "Enrichment engine status"},
                    "401": {"description": "Unauthorized"},
                    "503": {"description": "Enrichment engine unavailable"},
                },
            }
        },
        "/api/health/slos": {
            "get": {
                "summary": "SLO compliance",
                "description": "Service Level Objective compliance check across all monitored metrics.",
                "operationId": "healthSLOs",
                "tags": ["Health"],
                "security": [{"BearerAuth": []}],
                "responses": {
                    "200": {"description": "SLO compliance results"},
                    "401": {"description": "Unauthorized"},
                    "503": {"description": "Monitoring module unavailable"},
                },
            }
        },
        "/api/health/slo": {
            "get": {
                "summary": "Per-module SLO compliance",
                "description": "Per-module SLO compliance report for all 22+ modules with rolling 1-hour window metrics.",
                "operationId": "healthSLOModules",
                "tags": ["Health"],
                "responses": {
                    "200": {"description": "Per-module SLO compliance report"},
                    "503": {"description": "SLO config module unavailable"},
                },
            }
        },
        "/api/health/anomalies": {
            "get": {
                "summary": "Anomaly detection status",
                "description": "Statistical anomaly detection using rolling mean + 3 sigma thresholds.",
                "operationId": "healthAnomalies",
                "tags": ["Health"],
                "responses": {
                    "200": {"description": "Anomaly detection results"},
                    "503": {"description": "Anomaly detector module unavailable"},
                },
            }
        },
        "/api/health/eval": {
            "get": {
                "summary": "Eval scores",
                "description": "Run the full evaluation framework and return aggregate test scores.",
                "operationId": "healthEval",
                "tags": ["Health"],
                "security": [{"BearerAuth": []}],
                "responses": {
                    "200": {
                        "description": "Evaluation results with per-category scores"
                    },
                    "401": {"description": "Unauthorized"},
                    "503": {"description": "Eval framework unavailable"},
                },
            }
        },
        "/api/deploy/ready": {
            "get": {
                "summary": "Deploy readiness gate",
                "description": "Blue-green deployment readiness probe. Returns 200 only when the server is fully warmed up (KB loaded, caches primed).",
                "operationId": "deployReady",
                "tags": ["Health"],
                "responses": {
                    "200": {"description": "Instance is fully warmed and ready"},
                    "503": {"description": "Instance is still warming up"},
                },
            }
        },
        "/api/plan/negotiate": {
            "post": {
                "summary": "Multi-agent plan negotiation",
                "description": "Run a multi-agent negotiation where BudgetAgent, ChannelAgent, and AudienceAgent collaborate to produce an optimal media plan allocation.",
                "operationId": "planNegotiate",
                "tags": ["Plan"],
                "requestBody": {
                    "required": True,
                    "content": {
                        "application/json": {
                            "schema": {
                                "type": "object",
                                "required": ["budget", "job_title"],
                                "properties": {
                                    "budget": {"type": "number", "example": 50000},
                                    "job_title": {
                                        "type": "string",
                                        "example": "Senior Software Engineer",
                                    },
                                    "location": {
                                        "type": "string",
                                        "example": "San Francisco, CA",
                                    },
                                    "industry": {
                                        "type": "string",
                                        "example": "Technology",
                                    },
                                    "goals": {
                                        "type": "array",
                                        "items": {"type": "string"},
                                    },
                                    "max_rounds": {"type": "integer", "example": 5},
                                },
                            }
                        }
                    },
                },
                "responses": {
                    "200": {"description": "Negotiation result with final proposal"},
                    "400": {"description": "Invalid request parameters"},
                    "503": {"description": "Multi-agent module unavailable"},
                },
            }
        },
        "/api/jobs/{job_id}": {
            "get": {
                "summary": "Poll async job",
                "description": "Check status of an asynchronous generation job. "
                "When completed, returns the binary file content.",
                "operationId": "pollJob",
                "tags": ["Generation"],
                "parameters": [
                    {
                        "name": "job_id",
                        "in": "path",
                        "required": True,
                        "schema": {"type": "string"},
                        "description": "Job ID returned by async /api/generate",
                    }
                ],
                "responses": {
                    "200": {
                        "description": "Job status or completed binary content",
                        "content": {
                            "application/json": {
                                "schema": {
                                    "type": "object",
                                    "properties": {
                                        "job_id": {"type": "string"},
                                        "status": {
                                            "type": "string",
                                            "enum": [
                                                "processing",
                                                "completed",
                                                "failed",
                                            ],
                                        },
                                        "progress_pct": {"type": "integer"},
                                        "created": {
                                            "type": "string",
                                            "format": "date-time",
                                        },
                                        "elapsed_seconds": {"type": "number"},
                                    },
                                },
                            },
                            "application/zip": {
                                "schema": {"type": "string", "format": "binary"},
                            },
                        },
                    },
                    "404": {"description": "Job not found or expired"},
                },
            }
        },
        "/api/admin/keys": {
            "post": {
                "summary": "Manage API keys",
                "description": "Create, list, or revoke tiered API keys.",
                "operationId": "adminKeys",
                "tags": ["Admin"],
                "security": [{"BearerAuth": []}],
                "requestBody": {
                    "required": True,
                    "content": {
                        "application/json": {
                            "schema": {
                                "type": "object",
                                "required": ["action"],
                                "properties": {
                                    "action": {
                                        "type": "string",
                                        "enum": ["create", "list", "revoke"],
                                    },
                                    "tier": {
                                        "type": "string",
                                        "enum": ["free", "pro", "enterprise"],
                                    },
                                    "label": {
                                        "type": "string",
                                        "description": "Human-readable key label",
                                    },
                                    "key": {
                                        "type": "string",
                                        "description": "Key to revoke (for revoke action)",
                                    },
                                },
                            }
                        }
                    },
                },
                "responses": {
                    "200": {"description": "Key management result"},
                    "401": {"description": "Unauthorized"},
                },
            }
        },
        "/api/admin/usage": {
            "get": {
                "summary": "Per-key usage dashboard",
                "description": "View rate limit usage statistics for all API keys.",
                "operationId": "adminUsage",
                "tags": ["Admin"],
                "security": [{"BearerAuth": []}],
                "responses": {
                    "200": {"description": "Usage statistics per key"},
                    "401": {"description": "Unauthorized"},
                },
            }
        },
    },
    "components": {
        "securitySchemes": {
            "BearerAuth": {
                "type": "http",
                "scheme": "bearer",
                "description": "Admin API key passed as Bearer token. "
                "Set ADMIN_API_KEY env var on the server.",
            }
        }
    },
    "tags": [
        {"name": "Generation", "description": "Media plan generation (sync and async)"},
        {"name": "Chat", "description": "Nova AI recruitment intelligence chatbot"},
        {
            "name": "Health",
            "description": "Health checks, readiness probes, monitoring",
        },
        {
            "name": "Admin",
            "description": "Administrative endpoints (require Bearer token)",
        },
    ],
}

# ═══════════════════════════════════════════════════════════════════════════════
# SWAGGER UI HTML (self-contained, loads from CDN)
# ═══════════════════════════════════════════════════════════════════════════════
_SWAGGER_UI_HTML = """<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>AI Media Planner - API Docs</title>
<link rel="stylesheet" href="https://unpkg.com/swagger-ui-dist@5/swagger-ui.css">
<style>
body { margin: 0; padding: 0; }
#swagger-ui .topbar { display: none; }
</style>
</head>
<body>
<div id="swagger-ui"></div>
<script src="https://unpkg.com/swagger-ui-dist@5/swagger-ui-bundle.js"></script>
<script>
SwaggerUIBundle({
  url: "/api/docs/openapi.json",
  dom_id: "#swagger-ui",
  deepLinking: true,
  presets: [SwaggerUIBundle.presets.apis, SwaggerUIBundle.SwaggerUIStandalonePreset],
  layout: "BaseLayout",
});
</script>
</body>
</html>"""

# ═══════════════════════════════════════════════════════════════════════════════
# ═══════════════════════════════════════════════════════════════════════════════
# SHARED PLANS STORE (Plan Sharing / Collaborative Editing)
# ═══════════════════════════════════════════════════════════════════════════════
_shared_plans: dict[str, dict] = {}
_shared_plans_lock = threading.Lock()
_SHARED_PLANS_MAX = 1000
_SHARED_PLANS_TTL = 86400.0  # 24 hours

_plan_feedback: dict[str, list[dict]] = {}
_plan_feedback_lock = threading.Lock()
_plan_feedback_ts: dict[str, float] = {}  # track last-update time per key
_PLAN_FEEDBACK_MAX = 1000
_PLAN_FEEDBACK_TTL = 86400.0  # 24 hours

# SCORECARD STORE (in-memory, for shareable plan scorecards)
# Values are (html_string, created_timestamp) tuples
_scorecards: dict[str, tuple[str, float]] = {}
_scorecards_lock = threading.Lock()
_SCORECARDS_MAX = 500
_SCORECARDS_TTL = 43200.0  # 12 hours

# ASYNC GENERATION JOB STORE
# ═══════════════════════════════════════════════════════════════════════════════
_generation_jobs: dict = {}
_generation_jobs_lock = threading.Lock()
_GENERATION_JOB_EXPIRY_SECONDS = (
    24 * 60 * 60
)  # 24 hours (was 30 min -- Slack users download hours later)

# ═══════════════════════════════════════════════════════════════════════════════
# PLAN RESULTS STORE (in-memory, TTL 30min, for on-screen dashboard)
# ═══════════════════════════════════════════════════════════════════════════════
_plan_results_store: dict[str, dict] = {}
_plan_results_lock = threading.Lock()
_PLAN_RESULTS_TTL_SECONDS = 24 * 60 * 60  # 24 hours (was 30 min)


def _extract_plan_json(data: dict) -> dict:
    """Extract a JSON summary from generation data for the on-screen dashboard."""
    budget_alloc = data.get("_budget_allocation") or {}
    synthesized = data.get("_synthesized") or {}
    enriched = data.get("_enriched") or {}
    roles_raw = data.get("target_roles") or data.get("roles") or []
    locations_raw = data.get("locations") or []
    industry = data.get("industry_label") or data.get("industry") or ""

    channels: list[dict] = []
    ch_allocs = budget_alloc.get("channel_allocations") or {}
    meta = budget_alloc.get("metadata") or {}
    total_budget = float(meta.get("total_budget") or 0)

    for ch_name, ch_data in ch_allocs.items():
        if not isinstance(ch_data, dict):
            continue
        dollar_amt = float(ch_data.get("dollar_amount") or 0)
        pct = float(ch_data.get("percentage") or ch_data.get("pct") or 0)
        if pct == 0 and total_budget > 0:
            pct = round(dollar_amt / total_budget * 100, 1)
        cpc = ch_data.get("cpc") or ch_data.get("cpc_range") or 0
        cpa = ch_data.get("cpa") or ch_data.get("cpa_range") or 0
        clicks = int(ch_data.get("projected_clicks") or 0)
        reason = ch_data.get("reason") or ch_data.get("recommendation") or ""
        channels.append(
            {
                "name": ch_name.replace("_", " ").title(),
                "allocation_pct": round(pct, 1),
                "budget": round(dollar_amt, 2),
                "cpc_range": (
                    f"${cpc:.2f}" if isinstance(cpc, (int, float)) else str(cpc)
                ),
                "cpa_range": (
                    f"${cpa:.2f}" if isinstance(cpa, (int, float)) else str(cpa)
                ),
                "projected_clicks": clicks,
                "recommended_reason": str(reason),
            }
        )
    channels.sort(key=lambda c: c["budget"], reverse=True)

    by_channel = {c["name"]: c["budget"] for c in channels}
    by_category: dict[str, float] = {}
    for ch in channels:
        cat = "Other"
        name_lower = ch["name"].lower()
        if "programmatic" in name_lower or "dsp" in name_lower:
            cat = "Programmatic"
        elif "board" in name_lower:
            cat = "Job Boards"
        elif "social" in name_lower:
            cat = "Social Media"
        elif "brand" in name_lower:
            cat = "Employer Branding"
        by_category[cat] = by_category.get(cat, 0) + ch["budget"]

    hiring_difficulty = str(
        synthesized.get("hiring_difficulty")
        or enriched.get("hiring_difficulty")
        or "moderate"
    )
    competition = str(
        synthesized.get("competition_level")
        or enriched.get("market_competition")
        or "medium"
    )
    salary_range = str(
        synthesized.get("salary_range") or enriched.get("salary_range") or "N/A"
    )
    demand_trend = str(
        synthesized.get("demand_trend") or enriched.get("demand_trend") or "stable"
    )

    recs: list[str] = []
    if channels:
        recs.append(
            f"Lead with {channels[0]['name']} for highest ROI based on industry benchmarks."
        )
    if len(channels) > 1:
        recs.append(
            f"Diversify across {len(channels)} channels to reduce single-source risk."
        )
    avg_cpc_vals = [
        c["budget"] / max(c["projected_clicks"], 1)
        for c in channels
        if c["projected_clicks"] > 0
    ]
    if avg_cpc_vals:
        avg = sum(avg_cpc_vals) / len(avg_cpc_vals)
        recs.append(
            f"Average effective CPC is ${avg:.2f} -- monitor and optimize weekly."
        )
    recs.append(
        "Review channel performance after 2 weeks and reallocate underperforming budget."
    )
    recs.append("Consider A/B testing job ad creatives across top 2 channels.")

    # ── Gold Standard enrichments ──
    gold_standard = data.get("_gold_standard") or {}

    result = {
        "channels": channels,
        "budget_summary": {
            "total": total_budget,
            "by_channel": by_channel,
            "by_category": by_category,
        },
        "market_insights": {
            "hiring_difficulty": hiring_difficulty,
            "competition_level": competition,
            "salary_range": salary_range,
            "demand_trend": demand_trend,
        },
        "recommendations": recs,
        "metadata": {
            "roles": roles_raw if isinstance(roles_raw, list) else [roles_raw],
            "locations": (
                locations_raw if isinstance(locations_raw, list) else [locations_raw]
            ),
            "industry": industry,
            "total_budget": total_budget,
            "generated_at": datetime.datetime.utcnow().isoformat() + "Z",
            "gold_standard_gates": list(gold_standard.keys()),
        },
    }

    # Attach Gold Standard sections if present
    if gold_standard.get("city_level_data"):
        result["city_level_data"] = gold_standard["city_level_data"]
    if gold_standard.get("clearance_segmentation"):
        result["clearance_segmentation"] = gold_standard["clearance_segmentation"]
    if gold_standard.get("competitor_mapping"):
        result["competitor_mapping"] = gold_standard["competitor_mapping"]
    if gold_standard.get("difficulty_framework"):
        result["difficulty_framework"] = gold_standard["difficulty_framework"]
    if gold_standard.get("channel_strategy"):
        result["channel_strategy"] = gold_standard["channel_strategy"]
    if gold_standard.get("budget_tiers"):
        result["budget_tiers"] = gold_standard["budget_tiers"]
    if gold_standard.get("activation_calendar"):
        result["activation_calendar"] = gold_standard["activation_calendar"]

    return result


def _store_plan_result(plan_id: str, data: dict) -> None:
    """Store extracted plan JSON with TTL for dashboard retrieval."""
    try:
        plan_json = _extract_plan_json(data)
        with _plan_results_lock:
            _plan_results_store[plan_id] = {"data": plan_json, "created": time.time()}
        logger.info("Plan result stored: %s", plan_id)
    except Exception as e:
        logger.error("Failed to store plan result: %s", e, exc_info=True)


def _upload_plan_to_drive(file_bytes: bytes, filename: str, mime_type: str) -> str:
    """Upload plan ZIP to Google Drive for permanent download link.

    Returns:
        Direct download URL, or empty string if upload fails or Drive not configured.
    """
    import base64 as _b64

    creds_b64 = os.environ.get("GOOGLE_SLIDES_CREDENTIALS_B64") or ""
    if not creds_b64:
        return ""

    try:
        creds_json = _b64.b64decode(creds_b64).decode("utf-8")
        creds_dict = json.loads(creds_json)
    except Exception:
        return ""

    try:
        from googleapiclient.discovery import build as _gapi_build
        from googleapiclient.http import MediaInMemoryUpload
        from google.oauth2.service_account import Credentials as _GCreds
    except ImportError:
        return ""

    try:
        scopes = ["https://www.googleapis.com/auth/drive"]
        creds = _GCreds.from_service_account_info(creds_dict, scopes=scopes)
        drive = _gapi_build("drive", "v3", credentials=creds)

        # S46 security: sanitize filename before Drive upload
        safe_name = (
            re.sub(r"[^\w\s\-_.]", "_", filename)[:200] if filename else "plan.zip"
        )
        file_meta = {"name": safe_name, "mimeType": mime_type}
        media = MediaInMemoryUpload(file_bytes, mimetype=mime_type)
        uploaded = (
            drive.files()
            .create(body=file_meta, media_body=media, fields="id,webContentLink")
            .execute()
        )

        file_id = uploaded.get("id") or ""
        if not file_id:
            return ""

        # Make publicly downloadable
        # S46 security fix: restrict to joveo.com domain instead of public
        drive.permissions().create(
            fileId=file_id,
            body={"role": "reader", "type": "domain", "domain": "joveo.com"},
        ).execute()

        return f"https://drive.google.com/uc?export=download&id={file_id}"
    except Exception as exc:
        logger.warning("Drive upload failed: %s", exc)
        return ""


def _cleanup_generation_jobs():
    """Background thread to clean up expired async generation jobs."""
    while True:
        try:
            time.sleep(300)  # Every 5 minutes
            now = time.time()
            with _generation_jobs_lock:
                # Mark stale "processing" jobs as failed (stuck > 10 minutes)
                # S29 v2: raised from 3min to 10min for quality-first pipeline
                for jid, jdata in _generation_jobs.items():
                    if (
                        jdata.get("status") == "processing"
                        and (now - jdata.get("created", 0)) > 600
                    ):
                        jdata["status"] = "failed"
                        jdata["error"] = "Generation timed out after 10 minutes"
                        jdata["result_bytes"] = None
                        logger.warning(
                            "Marked stale job %s as failed (stuck in processing > 10min)",
                            jid,
                        )

                expired = [
                    jid
                    for jid, jdata in _generation_jobs.items()
                    if (now - jdata.get("created", 0)) > _GENERATION_JOB_EXPIRY_SECONDS
                    or (
                        jdata.get("status") in ("completed", "failed")
                        and (now - jdata.get("created", 0)) > 600
                    )
                ]
                for jid in expired:
                    _generation_jobs.pop(jid, None)
            if expired:
                logger.info("Cleaned up %d expired generation jobs", len(expired))
        except Exception as e:
            logger.warning("Generation job cleanup error: %s", e)


_job_cleanup_thread = threading.Thread(
    target=_cleanup_generation_jobs, daemon=True, name="job-cleanup"
)
_job_cleanup_thread.start()

# ═══════════════════════════════════════════════════════════════════════════════
# REAL-TIME CPC/CPA MONITORING (S46)
# Runs every 6 hours, compares live data against KB benchmarks
# ═══════════════════════════════════════════════════════════════════════════════
_cpc_alerts: list = []
_cpc_alerts_lock = threading.Lock()


def _monitor_cpc_changes() -> None:
    """Background thread: compare live channel data against KB benchmarks every 6h."""
    import time as _t

    _t.sleep(60)  # Wait 60s after startup for data to load
    while True:
        try:
            live_file = os.path.join(DATA_DIR, "channel_benchmarks_live.json")
            kb_file = os.path.join(
                DATA_DIR, "recruitment_benchmarks_comprehensive_2026.json"
            )

            if not os.path.exists(live_file) or not os.path.exists(kb_file):
                _t.sleep(21600)  # 6 hours
                continue

            with open(live_file, encoding="utf-8") as f:
                live_data = json.load(f)
            with open(kb_file, encoding="utf-8") as f:
                kb_data = json.load(f)

            # Extract KB CPC benchmarks
            kb_cpc: dict[str, float] = {}
            cpc_section = (
                kb_data.get("platform_benchmarks", {})
                .get("cost_per_click", {})
                .get("data", [])
            ) or []
            for item in cpc_section:
                if isinstance(item, dict):
                    platform = item.get("platform", "")
                    median = item.get("cpc_median_usd") or item.get("cpc_max_usd") or 0
                    if platform and median:
                        kb_cpc[platform.lower()] = float(median)

            # Compare live vs KB
            new_alerts: list[dict] = []
            for platform_key, platform_data in live_data.items():
                if not isinstance(platform_data, dict):
                    continue
                live_cpc = platform_data.get("avg_cpc") or platform_data.get("cpc") or 0
                if not live_cpc:
                    continue

                for kb_name, kb_val in kb_cpc.items():
                    if (
                        platform_key.lower() in kb_name
                        or kb_name in platform_key.lower()
                    ):
                        change_pct = (
                            ((float(live_cpc) - kb_val) / kb_val) * 100
                            if kb_val > 0
                            else 0
                        )
                        if abs(change_pct) >= 15:  # 15% threshold
                            new_alerts.append(
                                {
                                    "platform": platform_key,
                                    "metric": "CPC",
                                    "kb_value": round(kb_val, 2),
                                    "live_value": round(float(live_cpc), 2),
                                    "change_pct": round(change_pct, 1),
                                    "direction": "up" if change_pct > 0 else "down",
                                    "timestamp": _t.strftime(
                                        "%Y-%m-%dT%H:%M:%SZ", _t.gmtime()
                                    ),
                                    "severity": (
                                        "high" if abs(change_pct) >= 30 else "medium"
                                    ),
                                }
                            )
                        break

            with _cpc_alerts_lock:
                _cpc_alerts.clear()
                _cpc_alerts.extend(new_alerts)
            if new_alerts:
                logger.info("CPC monitor: %d alerts generated", len(new_alerts))
            else:
                logger.info("CPC monitor: no significant changes detected")

        except Exception as exc:
            logger.error("CPC monitor error: %s", exc, exc_info=True)

        _t.sleep(21600)  # 6 hours


threading.Thread(target=_monitor_cpc_changes, daemon=True, name="cpc-monitor").start()

# ═══════════════════════════════════════════════════════════════════════════════
# TIERED API KEY RATE LIMITING
# ═══════════════════════════════════════════════════════════════════════════════
API_KEY_TIERS = {
    "free": {"rpm": 5, "rpd": 50},
    "pro": {"rpm": 30, "rpd": 1000},
    "enterprise": {"rpm": 100, "rpd": 10000},
}

# In-memory store: key -> {tier, label, created, revoked, usage_minute: [...], usage_day: [...]}
_api_keys_store: dict = {}
_api_keys_lock = threading.Lock()

# ═══════════════════════════════════════════════════════════════════════════════
# MONITORING & OBSERVABILITY
# ═══════════════════════════════════════════════════════════════════════════════
try:
    from monitoring import (
        configure_logging,
        get_metrics,
        get_persistence,
        health_check_liveness,
        health_check_readiness,
        init_metrics_persistence,
        GracefulShutdown,
    )

    configure_logging(os.environ.get("LOG_LEVEL", "INFO"))
    _metrics = get_metrics()
    _shutdown = GracefulShutdown(timeout=30.0)
    init_metrics_persistence()
    logger.info("Monitoring module loaded successfully")
except ImportError as _mon_err:
    logger.warning("monitoring module not available: %s", _mon_err)
    _metrics = None
    _shutdown = None

    def health_check_liveness():
        return {
            "status": "ok",
            "version": "4.0.0",
            "timestamp": datetime.datetime.now().isoformat(),
        }

    def health_check_detailed():
        """Detailed health check with uptime, version, and subsystem checks."""
        return _build_health_response()

    def health_check_readiness():
        return {"status": "healthy", "version": "4.0.0", "metrics_persisted": False}

    def get_persistence():
        return None

    def init_metrics_persistence():
        pass


def _run_with_timeout(
    fn: callable, timeout_sec: float = 3.0, default: Any = None
) -> Any:
    """Run a callable in a daemon thread with a timeout.

    Returns the result of fn() if it completes within timeout_sec,
    otherwise returns default. Used by health checks to prevent any
    single external call from blocking the entire response.

    Args:
        fn: Zero-argument callable to execute.
        timeout_sec: Maximum seconds to wait.
        default: Value returned on timeout or error.

    Returns:
        Result of fn() or default.
    """
    result_container: list = []

    def _target() -> None:
        try:
            result_container.append(fn())
        except Exception:
            pass

    t = threading.Thread(target=_target, daemon=True)
    t.start()
    t.join(timeout=timeout_sec)
    if result_container:
        return result_container[0]
    return default


def _build_health_response() -> dict:
    """Build the detailed health check response with subsystem status.

    Centralised so both the except-block fallback and the
    'if not in dir()' guard share the same logic.

    Includes a 5-second time guard: if the check takes longer than 5s,
    partial results are returned with a warning so Render health probes
    never time out.  Individual external checks (Upstash, vector search,
    etc.) are time-boxed to 3 seconds each via _run_with_timeout.
    """
    _health_start = time.time()
    _HEALTH_TIMEOUT = 5.0  # seconds -- hard cap for the entire check

    kb_loaded = _knowledge_base is not None and len(_knowledge_base) > 0
    data_dir_exists = os.path.isdir(DATA_DIR)
    is_healthy = kb_loaded and data_dir_exists

    # Supabase connectivity check
    _sb_connected = False
    try:
        _sb_connected = bool(_supabase_data_available)
    except NameError:
        pass

    # Enrichment daemon check
    _enrich_running = False
    try:
        _enrich_running = bool(_data_enrichment_available)
    except NameError:
        pass

    # LLM router availability
    _llm_available = False
    try:
        _llm_available = _lazy_llm_router() is not None
    except Exception:
        pass

    # API integrations availability
    _api_integ = False
    try:
        _api_integ = bool(_api_integrations_available)
    except NameError:
        pass

    # New integrations availability
    _jobspy_ok = False
    _tavily_ok = False
    _vector_ok = False
    try:
        _jobspy_ok = bool(_jobspy_available)
    except NameError:
        pass
    try:
        _tavily_ok = bool(_tavily_available)
    except NameError:
        pass
    try:
        _vector_ok = bool(_vector_search_available)
    except NameError:
        pass

    # PostHog availability flag (key is NOT exposed in health response for security)
    _ph_configured = bool(
        (
            os.environ.get("POSTHOG_PROJECT_API_KEY")
            or os.environ.get("POSTHOG_API_KEY")
            or ""
        ).strip()
    )

    # Resolve metrics_persisted from the persistence singleton (with grace period)
    _metrics_ok = False
    try:
        _p = get_persistence()
        if _p is not None:
            _metrics_ok = _p.is_persisted or _p.is_within_startup_grace
    except Exception:
        pass

    # ---- Runtime metrics from MetricsCollector ----
    _rt_metrics: dict = {}
    try:
        if _metrics is not None and hasattr(_metrics, "get_metrics"):
            _mc_snap = _metrics.get_metrics()
            _total_req = _mc_snap.get("total_requests") or 0
            _total_err = _mc_snap.get("total_errors") or 0
            _err_rate = (_total_err / max(1, _total_req)) if _total_req else 0.0
            _latency_data = _mc_snap.get("latency_ms") or {}
            _total_client_err = _mc_snap.get("total_client_errors") or 0
            _rt_metrics = {
                "total_requests": _total_req,
                "total_errors": _total_err,
                "total_client_errors": _total_client_err,
                "error_rate": round(_err_rate, 4),
                "avg_latency_ms": _latency_data.get("p50", 0.0),
                "p99_latency_ms": _latency_data.get("p99", 0.0),
                "active_requests": _mc_snap.get("active_requests") or 0,
                "requests_per_minute": _mc_snap.get("requests_per_minute") or 0,
                "status_codes": _mc_snap.get("status_codes", {}),
            }
        else:
            _rt_metrics = {
                "total_requests": 0,
                "total_errors": 0,
                "error_rate": 0.0,
                "avg_latency_ms": 0.0,
                "p99_latency_ms": 0.0,
                "active_requests": 0,
                "requests_per_minute": 0,
            }
    except Exception:
        _rt_metrics = {
            "total_requests": 0,
            "total_errors": 0,
            "error_rate": 0.0,
            "avg_latency_ms": 0.0,
            "p99_latency_ms": 0.0,
            "active_requests": 0,
            "requests_per_minute": 0,
        }

    # ---- Module status for dashboard ----
    _modules: dict = {
        "knowledge_base": "ok" if kb_loaded else "down",
        "data_dir": "ok" if data_dir_exists else "down",
        "supabase": "ok" if _sb_connected else "down",
        "enrichment_daemon": "ok" if _enrich_running else "down",
        "llm_router": "ok" if _llm_available else "down",
        "api_integrations": "ok" if _api_integ else "down",
        "metrics_collector": "ok" if (_metrics is not None) else "down",
    }
    # Check Upstash Redis (L3 cache) connectivity -- time-boxed to 3s
    try:
        from upstash_cache import ping as _upstash_ping, _ENABLED as _upstash_on

        if _upstash_on:
            _ping_result = _run_with_timeout(
                _upstash_ping, timeout_sec=3.0, default=False
            )
            _modules["upstash_redis"] = "ok" if _ping_result else "timeout"
        else:
            _modules["upstash_redis"] = "not_configured"
    except ImportError:
        _modules["upstash_redis"] = "not_configured"

    # Check optional modules -- attempt import if not yet in sys.modules (self-healing)
    # Parallelized: all imports run concurrently with a shared 3s deadline
    # instead of sequential 3s x N = 12s+ worst case.
    _check_modules = ("resilience_router", "posthog_tracker", "nova", "eval_framework")
    _mod_already = {m: "ok" for m in _check_modules if sys.modules.get(m)}
    _mod_missing = [m for m in _check_modules if m not in _mod_already]
    _modules.update(_mod_already)

    if _mod_missing:
        from concurrent.futures import ThreadPoolExecutor as _HealthPool

        def _try_import(name: str) -> tuple[str, bool]:
            """Attempt to import a module, returning (name, success)."""
            try:
                __import__(name)
                return (name, True)
            except Exception:
                return (name, False)

        # S61 FIX: do NOT use `with _HealthPool(...) as _hp:`.  The
        # context-manager __exit__ blocks on shutdown(wait=True), so a
        # single slow import locks the health check (stress test measured
        # /api/health at 8.5s, burning 17% of each worker's capacity).
        # Use finally + shutdown(wait=False, cancel_futures=True).
        _hp = _HealthPool(max_workers=len(_mod_missing))
        try:
            _futs = {_hp.submit(_try_import, m): m for m in _mod_missing}
            for _fut in _futs:
                try:
                    _name, _ok = _fut.result(timeout=3.0)
                    _modules[_name] = "ok" if _ok else "down"
                    if _ok:
                        logger.info(
                            f"Self-healed module '{_name}' -- imported on health check"
                        )
                except TimeoutError:
                    _modules[_futs[_fut]] = "down"
        finally:
            try:
                _hp.shutdown(wait=False, cancel_futures=True)
            except TypeError:
                _hp.shutdown(wait=False)

    # ---- LLM provider info for dashboard ----
    _llm_info: dict = {}
    try:
        if _llm_available:
            _router = _lazy_llm_router()
            if _router and hasattr(_router, "get_provider_status"):
                _llm_info = _router.get_provider_status()
            elif _router and hasattr(_router, "providers"):
                for _pname in _router.providers:
                    _llm_info[_pname] = {"status": "ok", "available": True}
            else:
                # Fallback: just report the router is available
                _llm_info["llm_router"] = {"status": "ok", "available": True}
    except Exception:
        _llm_info["llm_router"] = {"status": "degraded", "available": False}

    result: dict = {
        "status": "healthy" if is_healthy else "unhealthy",
        "uptime_seconds": round(time.time() - _SERVER_START_TIME, 2),
        "version": _DEPLOY_VERSION,
        "git_commit": _DEPLOY_GIT_COMMIT,
        "instance_id": _DEPLOY_INSTANCE_ID,
        "warmup_complete": _DEPLOY_WARMUP_COMPLETE,
        "timestamp": datetime.datetime.now(datetime.timezone.utc).isoformat(),
        "checks": {
            "knowledge_base": kb_loaded,
            "data_dir": data_dir_exists,
        },
        "metrics": _rt_metrics,
        "modules": _modules,
        "llm_providers": _llm_info,
        "metrics_persisted": _metrics_ok,
        "supabase_connected": _sb_connected,
        "enrichment_daemon_running": _enrich_running,
        "llm_router_available": _llm_available,
        "api_integrations_available": _api_integ,
        "jobspy_available": _jobspy_ok,
        "tavily_available": _tavily_ok,
        "vector_search_available": _vector_ok,
        "posthog_configured": _ph_configured,
        "upstash_redis_connected": _modules.get("upstash_redis") == "ok",
        "slack_alerts_available": _slack_alerts_available,
        "calendar_sync_available": _calendar_sync_available,
        "chroma_rag_available": _chroma_rag_available,
    }

    # P1-15: nova_conversations row count diagnostic
    try:
        from nova_persistence import get_conversation_count

        _conv_health = _run_with_timeout(
            get_conversation_count, timeout_sec=3.0, default={"error": "timeout"}
        )
        result["nova_conversations"] = _conv_health
    except ImportError:
        result["nova_conversations"] = {"error": "module_unavailable"}
    except Exception as _conv_exc:
        result["nova_conversations"] = {"error": str(_conv_exc)}

    # ---- Time-guarded module stats (skip if >5s elapsed) ----
    def _health_expired() -> bool:
        """Return True if the health check has exceeded its time budget."""
        return (time.time() - _health_start) > _HEALTH_TIMEOUT

    if _health_expired():
        result["warning"] = "Health check truncated (>5s) -- skipped module stats"
        result["health_check_ms"] = round((time.time() - _health_start) * 1000, 1)
        return result

    # Parallel external status checks (was sequential 4x3s = 12s worst case)
    _status_checks: dict[str, tuple] = {}
    if _vector_search_available:
        try:
            from vector_search import get_status as _vs_status

            _status_checks["vector_search"] = (_vs_status, {"status": "timeout"})
        except ImportError:
            result["vector_search"] = {"status": "import_error"}
    else:
        result["vector_search"] = {"status": "not_available"}

    if _slack_alerts_available:
        _status_checks["slack_alerts"] = (_slack_alerts_status, {"status": "timeout"})
    else:
        result["slack_alerts"] = {"status": "not_available"}

    if _calendar_sync_available:
        _status_checks["calendar_sync"] = (_cal_status, {"status": "timeout"})
    else:
        result["calendar_sync"] = {"status": "not_available"}

    if _chroma_rag_available:
        _status_checks["chroma_rag"] = (_chroma_stats, {"status": "timeout"})
    else:
        result["chroma_rag"] = {"status": "not_available"}

    if _status_checks:
        from concurrent.futures import (
            ThreadPoolExecutor as _StatusPool,
            as_completed as _as_done,
        )

        # S61 FIX: same pattern -- context manager's shutdown(wait=True)
        # was making /api/health wait for all status checks even when the
        # 3.5s deadline elapsed.
        _sp = _StatusPool(max_workers=len(_status_checks))
        try:
            _sfuts = {
                _sp.submit(fn): (name, default)
                for name, (fn, default) in _status_checks.items()
            }
            try:
                for _sfut in _as_done(_sfuts, timeout=3.5):
                    _sname, _sdefault = _sfuts[_sfut]
                    try:
                        result[_sname] = _sfut.result(timeout=0.1)
                    except TimeoutError:
                        result[_sname] = _sdefault
                    except Exception as _status_err:
                        logger.warning(
                            f"Health status check '{_sname}' failed: {_status_err}"
                        )
                        result[_sname] = _sdefault
            except TimeoutError:
                # _as_done timed out -- mark remaining checks as default and move on
                for _fut2, (_sname2, _sdefault2) in _sfuts.items():
                    if _sname2 not in result:
                        result[_sname2] = _sdefault2
        finally:
            try:
                _sp.shutdown(wait=False, cancel_futures=True)
            except TypeError:
                _sp.shutdown(wait=False)

    if _health_expired():
        result["warning"] = (
            "Health check truncated (>5s) -- skipped module stats after parallel checks"
        )
        result["health_check_ms"] = round((time.time() - _health_start) * 1000, 1)
        return result

    # Request coalescing stats
    try:
        from request_coalescing import (
            get_stats as _coal_stats,
            _ENABLED as _coal_enabled,
        )

        if _coal_enabled:
            result["request_coalescing"] = _coal_stats()
        else:
            result["request_coalescing"] = {"status": "not_configured"}
    except ImportError:
        result["request_coalescing"] = {"status": "not_installed"}

    # Feature store status
    try:
        from feature_store import get_feature_store as _get_fs

        _fs = _get_fs()
        result["feature_store"] = _fs.get_all_features()
        _modules["feature_store"] = "ok" if _fs._initialized else "down"
    except ImportError:
        result["feature_store"] = {"initialized": False}
        _modules["feature_store"] = "not_installed"
    except Exception as _fs_health_err:
        logger.warning("Feature store health check failed: %s", _fs_health_err)
        result["feature_store"] = {"initialized": False, "error": str(_fs_health_err)}
        _modules["feature_store"] = "down"

    # Circuit breaker mesh status
    try:
        from circuit_breaker_mesh import get_circuit_mesh as _get_cb_mesh

        _cb_mesh = _get_cb_mesh()
        _mesh_status = _cb_mesh.get_mesh_status()
        result["circuit_breaker_mesh"] = {
            "total_providers": _mesh_status.get("total_providers", 0),
            "states": _mesh_status.get("states", {}),
            "avg_health_score": _mesh_status.get("avg_health_score", 0.0),
            "mesh_uptime_s": _mesh_status.get("mesh_uptime_s", 0.0),
        }
    except ImportError:
        result["circuit_breaker_mesh"] = {"status": "not_available"}
    except Exception as _cb_err:
        logger.warning(f"Circuit breaker mesh status failed: {_cb_err}")
        result["circuit_breaker_mesh"] = {"status": "error"}

    # Canvas engine stats
    try:
        from canvas_engine import get_canvas_stats as _canvas_stats_fn

        result["canvas_engine"] = _canvas_stats_fn()
    except ImportError:
        result["canvas_engine"] = {"status": "not_installed"}
    except Exception as _canvas_err:
        logger.warning("Canvas engine health check failed: %s", _canvas_err)
        result["canvas_engine"] = {"status": "error", "detail": str(_canvas_err)}

    # Benchmark file freshness check
    try:
        from data_enrichment import check_benchmark_freshness

        bench_fresh = check_benchmark_freshness()
        any_stale = any(v.get("stale", False) for v in bench_fresh.values())
        result["benchmark_freshness"] = bench_fresh
        if any_stale:
            result["benchmark_stale_warning"] = True
    except ImportError:
        pass
    except (ValueError, KeyError, TypeError, OSError) as bfe:
        logger.warning(f"Benchmark freshness check failed: {bfe}")

    if _health_expired():
        result["warning"] = "Health check truncated (>5s) -- partial module stats"
        result["health_check_ms"] = round((time.time() - _health_start) * 1000, 1)
        return result

    # Edge routing (geographic LLM provider routing) stats
    try:
        from edge_routing import get_edge_routing_stats as _get_edge_stats

        result["edge_routing"] = _get_edge_stats()
        _modules["edge_routing"] = "ok"
    except ImportError:
        result["edge_routing"] = {"status": "not_installed"}
        _modules["edge_routing"] = "not_installed"
    except Exception as _er_err:
        logger.warning(f"Edge routing health check failed: {_er_err}")
        result["edge_routing"] = {"status": "error"}
        _modules["edge_routing"] = "down"

    # Morning brief system status
    try:
        from morning_brief import get_brief_stats

        result["morning_brief"] = get_brief_stats()
    except ImportError:
        result["morning_brief"] = {"status": "not_installed"}
    except Exception as _mb_err:
        logger.warning(f"Morning brief health check failed: {_mb_err}")
        result["morning_brief"] = {"status": "error"}

    # Market pulse system status
    try:
        from market_pulse import get_pulse_stats

        result["market_pulse"] = get_pulse_stats()
        _modules["market_pulse"] = "ok"
    except ImportError:
        result["market_pulse"] = {"available": False, "status": "not_installed"}
        _modules["market_pulse"] = "not_installed"
    except Exception as _mp_err:
        logger.warning(f"Market pulse health check failed: {_mp_err}")
        result["market_pulse"] = {"available": False, "status": "error"}
        _modules["market_pulse"] = "down"

    # Role taxonomy stats
    try:
        from role_taxonomy import get_role_taxonomy as _get_rt

        result["role_taxonomy"] = _get_rt().get_taxonomy_stats()
    except ImportError:
        result["role_taxonomy"] = {"status": "not_installed"}
    except Exception as _rt_err:
        logger.warning(f"Role taxonomy health check failed: {_rt_err}")
        result["role_taxonomy"] = {"status": "error"}

    # Plan templates marketplace stats
    try:
        from plan_templates import get_templates_stats as _get_tmpl_stats

        result["plan_templates"] = _get_tmpl_stats()
    except ImportError:
        result["plan_templates"] = {"status": "not_installed"}
    except Exception as _tmpl_err:
        logger.warning(f"Plan templates health check failed: {_tmpl_err}")
        result["plan_templates"] = {"status": "error"}

    # Anonymized benchmarking network stats
    try:
        from benchmarking import get_benchmarking_stats as _get_bench_stats

        result["benchmarking_network"] = _get_bench_stats()
        _modules["benchmarking_network"] = (
            "ok" if result["benchmarking_network"].get("status") == "ok" else "down"
        )
    except ImportError:
        result["benchmarking_network"] = {"status": "not_installed"}
        _modules["benchmarking_network"] = "not_installed"
    except Exception as _bench_err:
        logger.warning(f"Benchmarking network health check failed: {_bench_err}")
        result["benchmarking_network"] = {"status": "error"}
        _modules["benchmarking_network"] = "down"

    # Adaptive rate limiter stats
    try:
        from rate_limiter_adaptive import get_rate_limiter_stats as _get_rl_stats

        result["rate_limiter"] = _get_rl_stats()
        _modules["rate_limiter"] = "ok"
    except ImportError:
        result["rate_limiter"] = {"status": "not_installed"}
        _modules["rate_limiter"] = "not_installed"
    except Exception as _rl_err:
        logger.warning(f"Rate limiter health check failed: {_rl_err}")
        result["rate_limiter"] = {"status": "error"}
        _modules["rate_limiter"] = "down"

    if _health_expired():
        result["warning"] = "Health check truncated (>5s) -- partial module stats"
        result["health_check_ms"] = round((time.time() - _health_start) * 1000, 1)
        return result

    # Outcome engine stats
    try:
        from outcome_engine import get_outcome_stats

        result["outcome_engine"] = get_outcome_stats()
        _modules["outcome_engine"] = "ok"
    except ImportError:
        result["outcome_engine"] = {"status": "not_installed"}
        _modules["outcome_engine"] = "not_installed"
    except Exception as _oe_err:
        logger.warning(f"Outcome engine health check failed: {_oe_err}")
        result["outcome_engine"] = {"status": "error"}
        _modules["outcome_engine"] = "down"

    # Outcome pipeline stats (closed-loop feedback)
    try:
        from outcome_pipeline import get_pipeline_stats as _get_pipeline_stats

        result["outcome_pipeline"] = _get_pipeline_stats()
        _modules["outcome_pipeline"] = "ok"
    except ImportError:
        result["outcome_pipeline"] = {"status": "not_installed"}
        _modules["outcome_pipeline"] = "not_installed"
    except Exception as _op_err:
        logger.warning(f"Outcome pipeline health check failed: {_op_err}")
        result["outcome_pipeline"] = {"status": "error"}
        _modules["outcome_pipeline"] = "down"

    # Attribution dashboard stats
    try:
        from attribution_dashboard import get_attribution_stats as _get_attr_stats

        result["attribution_dashboard"] = _get_attr_stats()
        _modules["attribution_dashboard"] = "ok"
    except ImportError:
        result["attribution_dashboard"] = {"status": "not_installed"}
        _modules["attribution_dashboard"] = "not_installed"
    except Exception as _attr_err:
        logger.warning(f"Attribution dashboard health check failed: {_attr_err}")
        result["attribution_dashboard"] = {"status": "error"}
        _modules["attribution_dashboard"] = "down"

    # Event store stats (legacy)
    try:
        from event_store import get_event_store_stats

        result["event_store"] = get_event_store_stats()
        _modules["event_store"] = "ok"
    except ImportError:
        result["event_store"] = {"status": "not_installed"}
        _modules["event_store"] = "not_installed"
    except Exception as _es_err:
        logger.warning(f"Event store health check failed: {_es_err}")
        result["event_store"] = {"status": "error"}
        _modules["event_store"] = "down"

    # Plan events state machine (event-sourced with fork/merge/undo)
    try:
        from plan_events import get_event_store_stats as _get_pe_stats

        result["plan_events"] = _get_pe_stats()
        _modules["plan_events"] = "ok"
    except ImportError:
        result["plan_events"] = {"status": "not_installed"}
        _modules["plan_events"] = "not_installed"
    except Exception as _pe_err:
        logger.warning(f"Plan events health check failed: {_pe_err}")
        result["plan_events"] = {"status": "error"}
        _modules["plan_events"] = "down"

    if _health_expired():
        result["warning"] = "Health check truncated (>5s) -- partial module stats"
        result["health_check_ms"] = round((time.time() - _health_start) * 1000, 1)
        return result

    # Market signals engine stats
    try:
        from market_signals import get_signal_stats as _get_signal_stats

        result["market_signals"] = _get_signal_stats()
        _modules["market_signals"] = (
            "ok" if result["market_signals"].get("status") == "ok" else "down"
        )
    except ImportError:
        result["market_signals"] = {"status": "not_installed"}
        _modules["market_signals"] = "not_installed"
    except Exception as _ms_err:
        logger.warning(f"Market signals health check failed: {_ms_err}")
        result["market_signals"] = {"status": "error"}
        _modules["market_signals"] = "down"

    # ATS widget stats
    try:
        from ats_widget import get_widget_stats

        result["ats_widget"] = get_widget_stats()
        _modules["ats_widget"] = "ok"
    except ImportError:
        result["ats_widget"] = {"status": "not_installed"}
        _modules["ats_widget"] = "not_installed"
    except Exception as _aw_err:
        logger.warning(f"ATS widget health check failed: {_aw_err}")
        result["ats_widget"] = {"status": "error"}
        _modules["ats_widget"] = "down"

    # Prediction model stats
    try:
        from prediction_model import get_prediction_stats

        result["prediction_model"] = get_prediction_stats()
        _modules["prediction_model"] = "ok"
    except ImportError:
        result["prediction_model"] = {"status": "not_installed"}
        _modules["prediction_model"] = "not_installed"
    except Exception as _pm_err:
        logger.warning(f"Prediction model health check failed: {_pm_err}")
        result["prediction_model"] = {"status": "error"}
        _modules["prediction_model"] = "down"

    # Embed widget stats (ATS integration tracking)
    with _embed_stats_lock:
        result["embed_widget"] = {
            "loads": _embed_stats["loads"],
            "js_serves": _embed_stats["js_serves"],
            "demo_views": _embed_stats["demo_views"],
            "unique_hosts": len(_embed_stats["hosts"]),
            "first_seen": _embed_stats["first_seen"],
            "last_seen": _embed_stats["last_seen"],
        }

    result["health_check_ms"] = round((time.time() - _health_start) * 1000, 1)
    return result


# Ensure health_check_detailed always exists (monitoring module may not export it)
if "health_check_detailed" not in dir():

    def health_check_detailed():
        """Detailed health check with uptime, version, and subsystem checks."""
        return _build_health_response()


# Data matrix health monitor (background checks every 12h)
try:
    from data_matrix_monitor import get_data_matrix_monitor

    _data_matrix = get_data_matrix_monitor()
    _data_matrix.start_background()
    logger.info("Data matrix monitor started (checks every 12h)")
except ImportError as _dm_err:
    logger.warning("data_matrix_monitor not available: %s", _dm_err)
    _data_matrix = None

# Autonomous QC engine (twice-daily tests). S57 audit: the "weekly
# self-upgrade" feature was never implemented -- log message and docstring
# corrected to reflect what the engine actually does (periodic health
# sampling). Adding real self-upgrade is a separate initiative.
try:
    from auto_qc import get_auto_qc

    _auto_qc = get_auto_qc()
    _auto_qc.start_background()
    logger.info("AutoQC engine started (health checks every 12h)")
except ImportError as _aqc_err:
    logger.warning("auto_qc not available: %s", _aqc_err)
    _auto_qc = None

# Monitoring-to-alerting bridge (60s SLO checks -> alerts)
try:
    from monitoring import start_alert_bridge as _start_alert_bridge

    _start_alert_bridge()
    logger.info("Monitoring alert bridge started")
except ImportError:
    logger.warning("monitoring alert bridge not available")

# Data enrichment engine (hourly freshness checks)
# Opt-in via ENABLE_DATA_ENRICHMENT=1 to avoid $8-15/mo of redundant external-API
# fan-out (BLS/FRED/Adzuna/Census/Firecrawl/etc.) that is already cached in the
# Supabase L3 layer. Manual/admin runs can still call start_enrichment() directly.
try:
    from data_enrichment import start_enrichment

    _enable_data_enrichment = (
        os.environ.get("ENABLE_DATA_ENRICHMENT") or ""
    ).strip().lower() in {"1", "true", "yes", "on"}
    if _enable_data_enrichment:
        start_enrichment()
        _data_enrichment_available = True
        logger.info("data_enrichment: ENABLED (hourly external API fan-out)")
    else:
        # Module is importable and usable for on-demand/admin runs, but the
        # background scheduler is not started. Flip ENABLE_DATA_ENRICHMENT=1
        # on Render to re-enable hourly refreshes.
        _data_enrichment_available = False
        logger.info(
            "data_enrichment: DISABLED (set ENABLE_DATA_ENRICHMENT=1 to enable)"
        )
except ImportError as _de_err:
    logger.warning("data_enrichment not available: %s", _de_err)
    _data_enrichment_available = False

# Email alert notifications (Resend API)
try:
    import email_alerts as _email_alerts

    if _email_alerts._is_enabled():
        logger.info("Email alerts enabled (Resend API)")
    else:
        logger.debug("Email alerts not configured (RESEND_API_KEY not set)")
except ImportError:
    _email_alerts = None
    logger.debug("email_alerts module not available")

# Supabase persistent cache (L3 cache layer)
try:
    from supabase_cache import start_cleanup_thread as _supa_start_cleanup

    _supa_start_cleanup(interval_hours=6)
    logger.info("Supabase cache cleanup thread started (every 6h)")
except ImportError:
    logger.debug("supabase_cache module not available")
except Exception as _supa_err:
    logger.debug("Supabase cache not configured: %s", _supa_err)

# ── Scheduled daily knowledge-base backup ─────────────────────────────────
_KB_BACKUP_INTERVAL = 86400  # 24 hours


def _scheduled_kb_backup() -> None:
    """Run a knowledge-base backup on a recurring timer.

    Uses threading.Timer so it re-schedules itself after each run,
    matching the existing pattern for background tasks.
    """
    try:
        from scripts.backup_kb import backup_knowledge_base

        result = backup_knowledge_base(DATA_DIR)
        logger.info("Scheduled KB backup succeeded: %s", Path(result).name)
    except Exception as bk_err:
        logger.error("Scheduled KB backup failed: %s", bk_err, exc_info=True)
    finally:
        # Re-schedule next run
        _timer = threading.Timer(_KB_BACKUP_INTERVAL, _scheduled_kb_backup)
        _timer.daemon = True
        _timer.name = "kb-backup"
        _timer.start()


if os.path.isdir(DATA_DIR):
    # S62 anti-reconvergence: first KB backup fires at a random offset in
    # [0, _KB_BACKUP_INTERVAL) so the 24h timer doesn't land on the same
    # wall-clock minute as other long timers (qdrant_ping 12h,
    # data_matrix_monitor 12h). secrets.SystemRandom avoids predictable jitter.
    _kb_backup_jitter = secrets.SystemRandom().uniform(0.0, float(_KB_BACKUP_INTERVAL))
    _kb_backup_timer = threading.Timer(_kb_backup_jitter, _scheduled_kb_backup)
    _kb_backup_timer.daemon = True
    _kb_backup_timer.name = "kb-backup"
    _kb_backup_timer.start()
    logger.info(
        "KB backup scheduler started (first run in %.1fh, jittered from 24h base)",
        _kb_backup_jitter / 3600.0,
    )
else:
    logger.warning("KB backup scheduler not started (data/ directory not found)")


# ── S32: Qdrant keep-alive ping (every 12 hours) ──────────────────────────
# Prevents free-tier Qdrant cluster from being suspended due to inactivity.
_QDRANT_PING_INTERVAL = 12 * 60 * 60  # 12 hours in seconds


def _scheduled_qdrant_ping() -> None:
    """Ping Qdrant cluster to prevent inactivity suspension."""
    try:
        _q_url = os.environ.get("QDRANT_URL") or ""
        _q_key = os.environ.get("QDRANT_API_KEY") or ""
        if _q_url and _q_key:
            import urllib.request

            _req = urllib.request.Request(
                f"{_q_url.rstrip('/')}/collections",
                headers={"api-key": _q_key, "Content-Type": "application/json"},
            )
            with urllib.request.urlopen(_req, timeout=10) as _resp:
                logger.info("Qdrant keep-alive ping: %s", _resp.status)
        else:
            logger.debug("Qdrant keep-alive skipped (no URL/key configured)")
    except Exception as _qp_err:
        logger.warning("Qdrant keep-alive ping failed: %s", _qp_err)
    finally:
        _qt = threading.Timer(_QDRANT_PING_INTERVAL, _scheduled_qdrant_ping)
        _qt.daemon = True
        _qt.name = "qdrant-ping"
        _qt.start()


if os.environ.get("QDRANT_URL") and os.environ.get("QDRANT_API_KEY"):
    # S62 anti-reconvergence: the first ping used to fire at T+60s deterministically,
    # then every 12h. With two workers plus the data_matrix_monitor (also 12h), all
    # three reconverged every 12h at the same instant. Jitter the first-ping delay
    # to a uniform random point in [60s, 60s + _QDRANT_PING_INTERVAL); once offset,
    # the 12h period keeps the threads staggered forever.
    _qdrant_first_ping_delay = 60.0 + secrets.SystemRandom().uniform(
        0.0, float(_QDRANT_PING_INTERVAL)
    )
    _qdrant_ping_timer = threading.Timer(
        _qdrant_first_ping_delay, _scheduled_qdrant_ping
    )
    _qdrant_ping_timer.daemon = True
    _qdrant_ping_timer.name = "qdrant-ping"
    _qdrant_ping_timer.start()
    logger.info(
        "Qdrant keep-alive scheduler started (first ping in %.1fh, jittered from 12h base)",
        _qdrant_first_ping_delay / 3600.0,
    )


# Simple in-memory rate limiter
from collections import defaultdict

# Bounded thread pool for async Slack event processing (max 4 concurrent)
_slack_executor = ThreadPoolExecutor(max_workers=4, thread_name_prefix="slack-event")
_rate_limit_store = defaultdict(list)
_rate_limit_lock = threading.Lock()
_RATE_LIMIT_WINDOW = 60  # seconds
_RATE_LIMIT_MAX = 30  # requests per window per IP (raised for chat UX)

# Per-user rate limit layer (on top of per-IP). Keyed by lowercase email.
# Higher cap than per-IP so a single heavy-testing engineer is not throttled,
# but an authenticated user who blows through cannot take down the whole
# office NAT egress IP bucket for their colleagues.
# Shares _rate_limit_lock for thread safety (no separate lock needed).
_user_rate_limit_store: dict = defaultdict(list)
_USER_RATE_LIMIT_MAX = int(
    os.environ.get("USER_RATE_LIMIT_MAX", "120")
)  # requests per _RATE_LIMIT_WINDOW per authenticated user

# Global rate limit for /api/chat to prevent distributed API cost abuse
_GLOBAL_CHAT_RATE_LIMIT_MAX = int(
    os.environ.get("GLOBAL_CHAT_RATE_LIMIT", "120")
)  # per minute across all IPs
_global_chat_timestamps: list = []
_global_chat_lock = threading.Lock()


# ═══════════════════════════════════════════════════════════════════════════════
# INPUT SANITIZATION
# ═══════════════════════════════════════════════════════════════════════════════

_re_sanitize = re  # re already imported at top of file

_HTML_TAG_RE = _re_sanitize.compile(r"<[^>]+>")
_SCRIPT_BLOCK_RE = _re_sanitize.compile(
    r"<script[\s>].*?</script>", _re_sanitize.IGNORECASE | _re_sanitize.DOTALL
)
_STYLE_BLOCK_RE = _re_sanitize.compile(
    r"<style[\s>].*?</style>", _re_sanitize.IGNORECASE | _re_sanitize.DOTALL
)
_EVENT_HANDLER_RE = _re_sanitize.compile(
    r'\bon\w+\s*=\s*["\'][^"\']*["\']', _re_sanitize.IGNORECASE
)
_MAX_CHAT_INPUT_LENGTH = 10000
_EMAIL_RE = re.compile(r"^[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}$")


def _sanitize_chat_input(text: Optional[str]) -> str:
    """Sanitize user chat input to prevent XSS while preserving intent.

    HTML-encodes angle brackets so user questions about HTML/scripts are
    preserved as readable text (e.g., '<script>' becomes '&lt;script&gt;')
    rather than being silently stripped, which would lose the question context.
    Also removes null bytes and event handler attributes, and truncates
    to 10000 characters.

    Args:
        text: Raw user input string (may be None).

    Returns:
        Sanitized string safe for rendering, with user intent preserved.
    """
    if not text:
        return ""
    # Ensure text is a string (handles int/bool/dict inputs from malformed requests)
    if not isinstance(text, str):
        text = str(text)
    # Strip null bytes and other control characters (prevent injection/log poisoning)
    text = text.replace("\x00", "")
    # Remove event handlers (e.g. onerror="...") -- these are always malicious
    text = _EVENT_HANDLER_RE.sub("", text)
    # HTML-encode angle brackets to prevent XSS while preserving user intent.
    # A user asking "What about <script> tags in job ads?" should still have
    # their question understood, not silently stripped to "What about  tags in job ads?"
    text = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
    # Truncate to max length
    if len(text) > _MAX_CHAT_INPUT_LENGTH:
        text = text[:_MAX_CHAT_INPUT_LENGTH]
    return text


def _dedup_repetitive_input(text: str, max_repeats: int = 3) -> str:
    """Collapse excessively repeated words/phrases to prevent LLM confusion.

    Detects when a user query contains the same word or short phrase repeated
    many times (e.g., "San Francisco " * 50) and deduplicates to unique words
    while preserving order.  Short inputs (< 100 chars or < 10 words) are
    returned unchanged.

    Args:
        text: Sanitized user input string.
        max_repeats: Threshold above which repetition triggers dedup.

    Returns:
        Deduplicated string, or the original if no excessive repetition found.
    """
    if not text or len(text) < 100:
        return text
    words = text.split()
    if len(words) < 10:
        return text
    # Check 1-3 word n-grams for excessive repetition
    from collections import Counter

    for ngram_size in (1, 2, 3):
        ngrams = [
            " ".join(words[i : i + ngram_size])
            for i in range(len(words) - ngram_size + 1)
        ]
        counts = Counter(ngrams)
        most_common_phrase, most_common_count = counts.most_common(1)[0]
        if most_common_count > max_repeats and most_common_count > len(words) // 3:
            # Extremely repetitive -- collapse to unique words preserving order
            seen: set[str] = set()
            unique_words: list[str] = []
            for w in words:
                key = w.lower().strip()
                if key not in seen:
                    seen.add(key)
                    unique_words.append(w)
            return " ".join(unique_words)
    return text


# ═══════════════════════════════════════════════════════════════════════════════
# ROUTE-AWARE SLIDING WINDOW RATE LIMITER
# ═══════════════════════════════════════════════════════════════════════════════


class RateLimiter:
    """Thread-safe sliding window rate limiter with per-route limits."""

    _MAX_IPS = 10000  # Hard cap on tracked IPs to prevent memory exhaustion

    def __init__(self):
        self._requests = {}  # ip -> list of timestamps
        self._lock = threading.Lock()

    def is_allowed(self, ip, max_requests=30, window_seconds=60):
        """Return True if request is allowed, False if rate limited."""
        now = time.time()
        with self._lock:
            if ip not in self._requests:
                # Evict oldest entries if we hit the cap (DDoS protection)
                if len(self._requests) >= self._MAX_IPS:
                    oldest_ip = min(
                        self._requests,
                        key=lambda k: (
                            max(self._requests[k]) if self._requests[k] else 0
                        ),
                    )
                    del self._requests[oldest_ip]
                self._requests[ip] = []
            # Remove expired timestamps
            self._requests[ip] = [
                t for t in self._requests[ip] if now - t < window_seconds
            ]
            if len(self._requests[ip]) >= max_requests:
                return False
            self._requests[ip].append(now)
            return True

    def cleanup(self):
        """Remove stale entries older than 5 minutes."""
        now = time.time()
        with self._lock:
            stale = [
                ip
                for ip, times in self._requests.items()
                if not times or now - max(times) > 300
            ]
            for ip in stale:
                del self._requests[ip]
        if stale:
            logger.debug("RateLimiter cleanup: removed %d stale IP entries", len(stale))


# Separate rate limiter instances for different route groups so that
# hitting the limit on /api/generate does not consume quota for other routes.
_rl_generate = RateLimiter()  # /api/generate -- 10 req/min (expensive)
_rl_chat = RateLimiter()  # /api/chat, /api/chat/stream -- 15 req/min (LLM calls)
_rl_llm_heavy = RateLimiter()  # LLM-heavy analysis endpoints -- 10 req/min
_rl_portal = RateLimiter()  # /api/portal/* -- 20 req/min
_rl_general = RateLimiter()  # all other /api/* POST routes -- 30 req/min
_rl_copilot = RateLimiter()  # /api/copilot/* -- 30 req/min (lightweight)

# ═══════════════════════════════════════════════════════════════════════════════
# CSRF: Cookie-based double-submit pattern (stateless)
# ═══════════════════════════════════════════════════════════════════════════════


_CSRF_TOKEN_TTL_SECONDS: int = 14400  # 4 hours expiry (extended for long sessions)


def _generate_csrf_token() -> str:
    """Generate a CSRF token with embedded expiry timestamp.

    Format: <random_hex>.<unix_timestamp>
    The timestamp enables server-side expiry validation.
    """
    random_part = secrets.token_hex(32)
    expiry = int(time.time()) + _CSRF_TOKEN_TTL_SECONDS
    return f"{random_part}.{expiry}"


def _build_csrf_cookie(token: str, *, secure: bool) -> str:
    """Build the Set-Cookie header value for the CSRF token.

    Args:
        token: The CSRF token string.
        secure: Whether to include the Secure flag (True for HTTPS).

    Returns:
        A fully-formed Set-Cookie header value.
    """
    parts = [
        f"csrf_token={token}",
        "Path=/",
        "SameSite=Strict",
    ]
    if secure:
        parts.append("Secure")
    return "; ".join(parts)


def _build_session_cookie(token: str, secure: bool = False) -> str:
    """Build a HttpOnly session cookie for IDOR prevention.

    This cookie is separate from csrf_token: the CSRF cookie must be readable
    by JavaScript (no HttpOnly) so it can be sent in the X-CSRF-Token header.
    The session cookie is HttpOnly and used purely server-side for job ownership.

    Args:
        token: The session token string (same value as CSRF token at creation time).
        secure: Whether to include the Secure flag (True for HTTPS).

    Returns:
        A fully-formed Set-Cookie header value.
    """
    parts = [
        f"nova_session={token}",
        "Path=/",
        "SameSite=Strict",
        "HttpOnly",
    ]
    if secure:
        parts.append("Secure")
    return "; ".join(parts)


def _parse_cookie_value(cookie_header: str, name: str) -> str:
    """Extract a named value from a Cookie header string.

    Args:
        cookie_header: The raw Cookie header (e.g. "a=1; csrf_token=abc; b=2").
        name: The cookie name to extract.

    Returns:
        The cookie value, or empty string if not found.
    """
    for pair in cookie_header.split(";"):
        pair = pair.strip()
        if pair.startswith(f"{name}="):
            return pair[len(name) + 1 :]
    return ""


def _validate_csrf_double_submit(cookie_token: str, header_token: str) -> bool:
    """Validate CSRF using the double-submit cookie pattern.

    Compares the token from the csrf_token cookie (set by server, sent
    automatically by browser) with the token from the X-CSRF-Token header
    (set by JavaScript). Uses constant-time comparison to prevent timing
    attacks. Also validates token expiry (Bug #15 fix).

    Args:
        cookie_token: Token read from the csrf_token cookie.
        header_token: Token read from the X-CSRF-Token request header.

    Returns:
        True if both tokens are non-empty, match, and have not expired.
    """
    import hmac as _hmac_mod

    if not cookie_token or not header_token:
        return False
    if not _hmac_mod.compare_digest(cookie_token, header_token):
        return False
    # Bug #15 fix: Check expiry timestamp embedded in token
    try:
        parts = cookie_token.rsplit(".", 1)
        if len(parts) == 2:
            expiry = int(parts[1])
            if time.time() > expiry:
                return False
        else:
            # Tokens without expiry are no longer accepted (legacy compat removed)
            logger.warning("CSRF token rejected: missing expiry timestamp")
            return False
    except (ValueError, TypeError):
        logger.warning("CSRF token rejected: invalid expiry format")
        return False
    return True


def _rate_limiter_cleanup_loop():
    """Background thread: run cleanup on all rate limiter instances every 5 minutes."""
    while True:
        time.sleep(300)
        try:
            _rl_generate.cleanup()
            _rl_chat.cleanup()
            _rl_llm_heavy.cleanup()
            _rl_portal.cleanup()
            _rl_general.cleanup()
        except Exception as exc:
            logger.warning("Rate limiter cleanup error: %s", exc)


_rl_cleanup_thread = threading.Thread(
    target=_rate_limiter_cleanup_loop, daemon=True, name="rate-limiter-cleanup"
)
_rl_cleanup_thread.start()

_ALLOWED_ORIGINS = {
    "http://localhost:10000",
    "http://localhost:5001",
    "http://127.0.0.1:10000",
    "http://127.0.0.1:5001",
    "https://media-plan-generator.onrender.com",
    # S48: Embed widget origins for cross-product CORS
    "https://cg-automation.onrender.com",
    "https://geoviz-3d.vercel.app",
    "https://nova.joveo.com",
    "https://geoviz.joveo.com",
}

# ═══════════════════════════════════════════════════════════════════════════════
# EMBED WIDGET STATS (thread-safe counters for ATS embed tracking)
# ═══════════════════════════════════════════════════════════════════════════════
_embed_stats_lock = threading.Lock()
_EMBED_HOSTS_MAX = 1000  # Cap host tracking to prevent unbounded growth
_embed_stats: dict[str, Any] = {
    "loads": 0,
    "js_serves": 0,
    "demo_views": 0,
    "hosts": {},  # hostname -> count
    "first_seen": None,
    "last_seen": None,
}


def _record_embed_event(event: str = "load", host: str = "") -> None:
    """Record an embed widget event (thread-safe)."""
    with _embed_stats_lock:
        now = datetime.datetime.utcnow().isoformat() + "Z"
        if event == "load":
            _embed_stats["loads"] += 1
        _embed_stats["last_seen"] = now
        if not _embed_stats["first_seen"]:
            _embed_stats["first_seen"] = now
        if host:
            hosts: dict[str, int] = _embed_stats["hosts"]
            if host in hosts:
                hosts[host] += 1
            elif len(hosts) < _EMBED_HOSTS_MAX:
                hosts[host] = 1
            else:
                # At cap -- evict least-recently-seen host (lowest count as proxy)
                min_host = min(hosts, key=hosts.get)  # type: ignore[arg-type]
                del hosts[min_host]
                hosts[host] = 1


# ═══════════════════════════════════════════════════════════════════════════════
# NOVA CHAT STREAMING INFRASTRUCTURE (SSE)
# ═══════════════════════════════════════════════════════════════════════════════

# Active streaming sessions: conversation_id -> threading.Event (set = cancel)
_active_streams: dict[str, threading.Event] = {}
_active_streams_lock = threading.Lock()


def _register_stream(conversation_id: str) -> threading.Event:
    """Register a streaming session and return its cancel event.

    Args:
        conversation_id: Unique conversation identifier.

    Returns:
        A threading.Event that, when set, signals the stream should stop.
    """
    cancel_event = threading.Event()
    with _active_streams_lock:
        # Cancel any existing stream for this conversation
        existing = _active_streams.get(conversation_id)
        if existing:
            existing.set()
        _active_streams[conversation_id] = cancel_event
    return cancel_event


def _unregister_stream(conversation_id: str) -> None:
    """Remove a streaming session from the active registry.

    Args:
        conversation_id: Unique conversation identifier.
    """
    with _active_streams_lock:
        _active_streams.pop(conversation_id, None)


def _cancel_stream(conversation_id: str) -> bool:
    """Signal cancellation for an active stream.

    Args:
        conversation_id: Unique conversation identifier.

    Returns:
        True if a stream was found and cancelled, False otherwise.
    """
    with _active_streams_lock:
        event = _active_streams.get(conversation_id)
        if event:
            event.set()
            return True
    return False


# ═══════════════════════════════════════════════════════════════════════════════
# FILE ATTACHMENT PARSING FOR CHAT
# ═══════════════════════════════════════════════════════════════════════════════

# ---------------------------------------------------------------------------
# File upload validation constants
# ---------------------------------------------------------------------------
_ALLOWED_FILE_EXTENSIONS: set[str] = {
    "pdf",
    "csv",
    "xlsx",
    "xls",
    "txt",
    "json",
    "png",
    "jpg",
    "jpeg",
    "gif",
    "doc",
    "docx",
}

_BLOCKED_FILE_EXTENSIONS: set[str] = {
    "exe",
    "bat",
    "sh",
    "py",
    "js",
    "cmd",
    "com",
    "msi",
    "scr",
    "vbs",
    "vbe",
    "wsf",
    "wsh",
    "ps1",
    "jar",
    "php",
    "rb",
    "pl",
}


def _sanitize_filename(raw_name: str) -> str:
    """Sanitize a filename by stripping path traversal and HTML-encoding.

    Removes directory separators, parent directory references (../),
    null bytes, and HTML-encodes special characters to prevent XSS.

    Args:
        raw_name: Raw filename from upload.

    Returns:
        Sanitized filename safe for logging and display.
    """
    import html
    from pathlib import PurePosixPath

    if not raw_name:
        return "unnamed_file"

    # Strip path components (handles Unix / and Windows \\)
    clean = PurePosixPath(raw_name.replace("\\", "/")).name

    # Remove null bytes and control characters
    clean = "".join(c for c in clean if c.isprintable() and c != "\x00")

    # HTML-encode to prevent XSS via filenames
    clean = html.escape(clean, quote=True)

    return clean or "unnamed_file"


def _validate_file_extension(filename: str) -> tuple[bool, str]:
    """Validate file extension against whitelist and blocklist.

    Args:
        filename: Sanitized filename.

    Returns:
        Tuple of (is_valid, error_message). error_message is empty if valid.
    """
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    if not ext:
        return False, "File has no extension"

    if ext in _BLOCKED_FILE_EXTENSIONS:
        return False, f"Executable file type .{ext} is not allowed"

    if ext not in _ALLOWED_FILE_EXTENSIONS:
        return False, (
            f"File type .{ext} is not supported. "
            f"Allowed: {', '.join(sorted(_ALLOWED_FILE_EXTENSIONS))}"
        )

    return True, ""


def _parse_file_attachment(file_attachment: dict) -> str:
    """Parse a base64-encoded file attachment and extract text content.

    Validates filename (sanitization, extension whitelist, size) before processing.

    Args:
        file_attachment: Dict with 'name', 'type', 'data' (base64).

    Returns:
        Extracted text content, or empty string on failure.
    """
    import base64
    import csv

    # io is imported at module level (line 6)

    raw_name = file_attachment.get("name") or ""
    name = _sanitize_filename(raw_name)
    data_b64 = file_attachment.get("data") or ""
    if not data_b64:
        return ""

    # Validate file extension before any processing
    is_valid, err_msg = _validate_file_extension(name)
    if not is_valid:
        logger.warning("File attachment rejected (%s): %s", name, err_msg)
        return f"[Rejected: {err_msg}]"

    try:
        raw = base64.b64decode(data_b64)
    except (ValueError, TypeError) as exc:
        logger.error("File decode error for %s: %s", name, exc, exc_info=True)
        return ""

    ext = name.rsplit(".", 1)[-1].lower() if "." in name else ""

    # Enforce 5MB limit
    if len(raw) > 5 * 1024 * 1024:
        return "[File too large -- max 5MB]"

    try:
        if ext == "txt":
            return raw.decode("utf-8", errors="replace")[:20000]

        if ext == "csv":
            text = raw.decode("utf-8", errors="replace")
            reader = csv.reader(io.StringIO(text))
            rows = []
            for i, row in enumerate(reader):
                if i > 200:
                    rows.append(f"... (truncated, {i}+ rows total)")
                    break
                rows.append(",".join(row))
            return "\n".join(rows)

        if ext == "json":
            text = raw.decode("utf-8", errors="replace")
            # Pretty-print if valid JSON, truncate to 20k chars
            try:
                parsed = json.loads(text)
                return json.dumps(parsed, indent=2)[:20000]
            except json.JSONDecodeError:
                return text[:20000]

        if ext == "pdf":
            # Best-effort PDF text extraction (stdlib only)
            text = raw.decode("latin-1", errors="replace")
            # Extract text between stream/endstream markers
            import re

            chunks: list[str] = []
            for match in re.finditer(
                rb"stream\r?\n(.*?)\r?\nendstream", raw, re.DOTALL
            ):
                chunk = match.group(1)
                # Try to decode as text (works for uncompressed PDFs)
                try:
                    decoded = chunk.decode("utf-8", errors="ignore")
                    # Filter out binary noise
                    printable = "".join(
                        c for c in decoded if c.isprintable() or c in "\n\t "
                    )
                    if len(printable) > 20:
                        chunks.append(printable)
                except (UnicodeDecodeError, ValueError):
                    pass
            if chunks:
                return "\n".join(chunks)[:20000]
            return "[PDF content could not be extracted with stdlib -- upload as TXT for best results]"

        if ext == "xlsx":
            # Best-effort XLSX extraction (it's a zip of XML files)
            import zipfile
            import xml.etree.ElementTree as ET

            try:
                zf = zipfile.ZipFile(io.BytesIO(raw))
                # Read shared strings
                shared_strings: list[str] = []
                if "xl/sharedStrings.xml" in zf.namelist():
                    ss_xml = zf.read("xl/sharedStrings.xml")
                    root = ET.fromstring(ss_xml)
                    ns = "{http://schemas.openxmlformats.org/spreadsheetml/2006/main}"
                    for si in root.findall(f".//{ns}si"):
                        texts = [t.text or "" for t in si.findall(f".//{ns}t")]
                        shared_strings.append("".join(texts))

                # Read first sheet
                sheet_names = [
                    n for n in zf.namelist() if n.startswith("xl/worksheets/sheet")
                ]
                if not sheet_names:
                    return "[XLSX: no worksheets found]"
                sheet_xml = zf.read(sheet_names[0])
                root = ET.fromstring(sheet_xml)
                ns = "{http://schemas.openxmlformats.org/spreadsheetml/2006/main}"
                rows_out: list[str] = []
                for row_el in root.findall(f".//{ns}row"):
                    cells: list[str] = []
                    for cell in row_el.findall(f"{ns}c"):
                        val_el = cell.find(f"{ns}v")
                        val = val_el.text if val_el is not None else ""
                        # Check if shared string reference
                        if cell.get("t") == "s" and val:
                            idx = int(val)
                            val = (
                                shared_strings[idx]
                                if idx < len(shared_strings)
                                else val
                            )
                        cells.append(val or "")
                    rows_out.append(",".join(cells))
                    if len(rows_out) > 200:
                        rows_out.append("... (truncated)")
                        break
                return "\n".join(rows_out)
            except (zipfile.BadZipFile, ET.ParseError, KeyError) as exc:
                logger.error("XLSX parse error for %s: %s", name, exc, exc_info=True)
                return "[XLSX parsing failed]"

        return f"[Unsupported file type: {ext}]"

    except Exception as exc:
        logger.error("File attachment parse error for %s: %s", name, exc, exc_info=True)
        return f"[Error parsing file: {name}]"


# ═══════════════════════════════════════════════════════════════════════════════
# SHARED CHAT CONTEXT ENRICHMENT (C-03 dedup + C-02 parallel)
# ═══════════════════════════════════════════════════════════════════════════════


def _enrich_chat_context(data: dict, message: str) -> dict:
    """Enrich chat request *data* dict with file-attachment text and live API data.

    Parses any attached file, then runs up to 6 API enrichment calls
    **concurrently** (ThreadPoolExecutor, 8-second combined timeout) so that
    one slow or failing service does not block the others.

    Args:
        data: Mutable request dict -- ``message`` may be extended with
              attachment text and ``_api_context`` is set with enrichment
              results.
        message: The raw user message (before lowercase / truncation).

    Returns:
        The (possibly mutated) *data* dict for convenience.
    """
    from concurrent.futures import ThreadPoolExecutor

    # ── 1. File attachment parsing ──
    _file_att = data.get("file_attachment")
    if _file_att and isinstance(_file_att, dict):
        try:
            _file_text = _parse_file_attachment(_file_att)
            if _file_text:
                _fname = _file_att.get("name") or "file"
                data["message"] = (
                    data.get("message") or ""
                ) + f"\n\n[Attached file: {_fname}]\n{_file_text}"
                logger.info(
                    "Parsed file attachment %s (%d chars)",
                    _fname,
                    len(_file_text),
                )
        except (ValueError, TypeError, OSError, KeyError) as _fpe:
            logger.error("File attachment parse failed: %s", _fpe, exc_info=True)

    # ── 2. Parallel API enrichment ──
    if (
        not _api_integrations_available
        and not _jobspy_available
        and not _tavily_available
        and not (_vector_search_available and _vector_search)
    ):
        return data

    _chat_msg = (data.get("message") or "").lower()
    _chat_role = data.get("role") or data.get("context", {}).get("role") or ""
    _chat_api_ctx: dict = {}
    _ctx_lock = threading.Lock()

    # -- Helper: thread-safe merge into _chat_api_ctx --
    def _merge(key: str, value: object) -> None:
        if value:
            with _ctx_lock:
                _chat_api_ctx[key] = value

    # -- Individual enrichment callables (each with its own try/except) --

    def _enrich_onet() -> None:
        if not (_api_integrations_available and _api_onet):
            return
        if not any(
            kw in _chat_msg
            for kw in ("job", "role", "occupation", "career", "position")
        ):
            return
        try:
            _search_term = _chat_role or _chat_msg[:50]
            results = _api_onet.search_occupations(_search_term)
            _merge("onet_occupations", results[:5] if results else None)
            if results:
                logger.info("Enriched chat with onet occupation data")
        except urllib.error.HTTPError as exc:
            if exc.code in (401, 403):
                logger.warning(
                    "O*NET enrichment auth failed (HTTP %s) -- check ONET_API_KEY",
                    exc.code,
                )
            else:
                logger.warning("O*NET enrichment failed: HTTP %s", exc.code)
        except (urllib.error.URLError, OSError, ValueError, TypeError) as exc:
            logger.warning("O*NET enrichment failed: %s", exc)

    def _enrich_adzuna() -> None:
        if not (_api_integrations_available and _api_adzuna):
            return
        if not any(
            kw in _chat_msg
            for kw in ("salary", "pay", "compensation", "wage", "earning")
        ):
            return
        try:
            _sal_role = _chat_role or "software engineer"
            hist = _api_adzuna.get_salary_histogram(_sal_role, "us")
            _merge("adzuna_salary_histogram", hist)
            if hist:
                logger.info("Enriched chat with adzuna salary data")
        except (urllib.error.URLError, OSError, ValueError, TypeError) as exc:
            logger.error("Adzuna enrichment failed: %s", exc, exc_info=True)

    def _enrich_fred() -> None:
        if not (_api_integrations_available and _api_fred):
            return
        if not any(
            kw in _chat_msg
            for kw in (
                "market",
                "economy",
                "unemployment",
                "inflation",
                "economic",
                "recession",
            )
        ):
            return
        try:
            unemp = _api_fred.get_unemployment_rate()
            _merge("fred_unemployment", unemp)
            cpi = _api_fred.get_cpi_data(months=6)
            _merge("fred_cpi", cpi)
            if unemp or cpi:
                logger.info("Enriched chat with fred economic data")
        except (urllib.error.URLError, OSError, ValueError, TypeError) as exc:
            logger.error("FRED enrichment failed: %s", exc, exc_info=True)

    def _enrich_jobspy() -> None:
        if not (_jobspy_available and _jobspy_market_stats):
            return
        if not any(
            kw in _chat_msg
            for kw in (
                "job",
                "hire",
                "hiring",
                "recruit",
                "posting",
                "vacancy",
                "position",
                "openings",
                "talent",
            )
        ):
            return
        try:
            _js_role = _chat_role or _chat_msg[:40]
            stats = _jobspy_market_stats(str(_js_role))
            _merge("jobspy_market_stats", stats)
            if stats:
                logger.info("Enriched chat with jobspy market stats")
        except (ValueError, TypeError, KeyError, OSError) as exc:
            logger.error("JobSpy enrichment failed: %s", exc, exc_info=True)

    def _enrich_tavily() -> None:
        if not (_tavily_available and _tavily_search):
            return
        if not any(
            kw in _chat_msg
            for kw in (
                "news",
                "trend",
                "latest",
                "recent",
                "update",
                "what's happening",
                "industry",
            )
        ):
            return
        try:
            results = _tavily_search(f"recruitment hiring {_chat_msg[:60]}")
            _merge("tavily_web_results", results)
            if results:
                logger.info(
                    "Enriched chat with tavily search (%d results)", len(results)
                )
        except (urllib.error.URLError, OSError, ValueError, TypeError) as exc:
            logger.error("Tavily enrichment failed: %s", exc, exc_info=True)

    def _enrich_vector() -> None:
        if not (_vector_search_available and _vector_search):
            return
        try:
            results = _vector_search(_chat_msg, top_k=3)
            _merge("vector_kb_results", results)
            if results:
                logger.info(
                    "Enriched chat with vector search (%d KB matches)", len(results)
                )
        except (ValueError, TypeError, OSError) as exc:
            logger.error("Vector search enrichment failed: %s", exc, exc_info=True)

    # -- Dispatch all enrichment tasks concurrently (8s combined timeout) --
    enrichment_fns = [
        _enrich_onet,
        _enrich_adzuna,
        _enrich_fred,
        _enrich_jobspy,
        _enrich_tavily,
        _enrich_vector,
    ]

    # S61 FIX: do NOT use `with ThreadPoolExecutor(...) as pool:` here.
    # The context-manager __exit__ calls shutdown(wait=True), which blocks
    # until ALL submitted tasks finish -- exactly what caused the typo-query
    # hang (claudedocs/s56_typo_hang_rca.md). A single slow enrichment
    # (Tavily throttle, JobSpy scrape, Adzuna 502) would leave the SSE
    # handler unable to yield its first byte for up to 60s, and the
    # per-task `fut.result(timeout=8)` did nothing to unblock it.
    # Now we shutdown(wait=False, cancel_futures=True) so slow tasks are
    # abandoned after the 8s per-task budget.
    pool = ThreadPoolExecutor(max_workers=4, thread_name_prefix="chat-enrich")
    try:
        futures = [pool.submit(fn) for fn in enrichment_fns]
        for fut in futures:
            try:
                fut.result(timeout=8)
            except TimeoutError:
                logger.warning("Chat enrichment task timed out (8s)")
            except Exception as exc:
                logger.error("Chat enrichment task error: %s", exc, exc_info=True)
    except Exception as exc:
        logger.error(
            "ThreadPoolExecutor failed for chat enrichment: %s", exc, exc_info=True
        )
    finally:
        # cancel_futures available in Python 3.9+. Fall back gracefully if
        # running on an older interpreter.
        try:
            pool.shutdown(wait=False, cancel_futures=True)
        except TypeError:
            pool.shutdown(wait=False)

    if _chat_api_ctx:
        data["_api_context"] = _chat_api_ctx

    return data


def _generate_followup_questions(response_text: str, user_message: str) -> list[str]:
    """Generate 2-3 suggested follow-up questions based on the response.

    Args:
        response_text: The Nova response text.
        user_message: The original user question.

    Returns:
        List of 2-3 short follow-up question strings.
    """
    # S47: Enhanced with entity extraction for contextual, specific follow-ups
    import re as _re

    # S47: Skip follow-ups for greetings and trivial messages
    _msg_stripped = (user_message or "").strip().lower().rstrip("!.,?")
    _greeting_words = {
        "hi",
        "hello",
        "hey",
        "hola",
        "howdy",
        "sup",
        "yo",
        "wassup",
        "wazzup",
        "whaddup",
        "namaste",
        "greetings",
        "thanks",
        "thank you",
        "thx",
        "ty",
        "bye",
        "goodbye",
        "ok",
        "okay",
        "yes",
        "no",
        "sure",
    }
    if _msg_stripped in _greeting_words:
        return []
    # Also skip for "how are you" type messages
    _smalltalk_pats = [
        r"^how are you",
        r"^how\'?s it going",
        r"^what\'?s up",
        r"^who are you",
        r"^what can you do",
        r"^good (morning|afternoon|evening)",
    ]
    if any(_re.search(p, _msg_stripped) for p in _smalltalk_pats):
        return []

    followups: list[str] = []
    resp_lower = response_text.lower()
    msg_lower = user_message.lower()

    # ── Extract key entities from response for contextual questions ──
    _role_patterns = _re.findall(
        r"(?:registered nurse|software engineer|data scientist|product manager|"
        r"account executive|sales representative|marketing manager|project manager|"
        r"business analyst|financial analyst|mechanical engineer|civil engineer|"
        r"nursing|developer|designer|recruiter|accountant|pharmacist|physician|"
        r"electrician|technician|administrator|coordinator|specialist|analyst|"
        r"engineer|manager|director|executive|consultant)",
        resp_lower,
    )
    _mentioned_role = _role_patterns[0].title() if _role_patterns else ""

    _city_patterns = _re.findall(
        r"(?:new york|los angeles|chicago|houston|dallas|san francisco|seattle|"
        r"boston|atlanta|denver|phoenix|miami|austin|portland|minneapolis|"
        r"charlotte|tampa|san diego|philadelphia|washington|detroit|orlando|"
        r"nashville|las vegas|salt lake city|raleigh|pittsburgh|columbus|"
        r"indianapolis|san antonio|san jose|jacksonville|memphis|baltimore|"
        r"milwaukee|kansas city|sacramento|st\.? louis|cincinnati|cleveland)",
        resp_lower,
    )
    _mentioned_city = _city_patterns[0].title() if _city_patterns else ""

    # ── Topic-specific contextual follow-ups ──
    if any(kw in resp_lower for kw in ("salary", "compensation", "pay range", "wage")):
        if _mentioned_role and "compare" not in msg_lower:
            followups.append(
                f"How does {_mentioned_role} salary compare to similar roles?"
            )
        if _mentioned_city:
            followups.append(
                f"What is the cost of living adjustment for {_mentioned_city}?"
            )
        elif "national" not in msg_lower:
            followups.append("How does this compare to national averages?")
        if "trend" not in msg_lower:
            followups.append("What are the salary trends over the past year?")
        if "remote" not in msg_lower:
            followups.append("How do remote salaries differ for this role?")

    elif any(
        kw in resp_lower
        for kw in ("cpc", "cpa", "cost per click", "cost per application")
    ):
        if _mentioned_role:
            followups.append(
                f"What channels give the best CPA for {_mentioned_role} roles?"
            )
        if "optimize" not in msg_lower:
            followups.append("How can I reduce these costs while maintaining quality?")
        if "trend" not in msg_lower:
            followups.append("How have these costs changed over the past quarter?")

    elif any(
        kw in resp_lower for kw in ("indeed", "linkedin", "job board", "publisher")
    ):
        if _mentioned_role:
            followups.append(
                f"Which niche job boards work best for {_mentioned_role} roles?"
            )
        elif "cost" not in msg_lower and "cpc" not in msg_lower:
            followups.append("What are the typical CPC benchmarks for these platforms?")
        if "roi" not in msg_lower and "best" not in msg_lower:
            followups.append("Which platform gives the best ROI for this type of role?")

    elif any(
        kw in resp_lower for kw in ("budget", "allocation", "spend", "media plan")
    ):
        if "optimize" not in msg_lower:
            followups.append(
                "How can I optimize this budget allocation for better ROI?"
            )
        followups.append("What KPIs should I track to measure this campaign's success?")

    elif any(
        kw in resp_lower for kw in ("hiring difficulty", "demand", "supply", "shortage")
    ):
        if _mentioned_city:
            followups.append(
                f"What sourcing strategies work best in {_mentioned_city}?"
            )
        else:
            followups.append("What sourcing strategies work best in this market?")
        if "remote" not in msg_lower:
            followups.append("How does remote hiring affect these numbers?")
        if "competitor" not in msg_lower:
            followups.append("Who are the top employers competing for this talent?")

    elif any(
        kw in resp_lower for kw in ("compliance", "regulation", "legal", "requirement")
    ):
        followups.append(
            "What are the most common compliance violations in job postings?"
        )
        followups.append("How do these requirements differ by state?")

    # General recruitment fallback
    if not followups:
        if any(kw in resp_lower for kw in ("recruit", "hire", "talent", "candidate")):
            if _mentioned_role:
                followups.append(
                    f"What are the best practices for hiring {_mentioned_role} talent?"
                )
            else:
                followups.append(
                    "What are the best practices for this hiring scenario?"
                )
            followups.append("Can you break this down by region?")
        else:
            followups.append("Can you provide more details on this topic?")
            if _mentioned_role:
                followups.append(
                    f"What salary and market data is available for {_mentioned_role} roles?"
                )
            else:
                followups.append("What recruitment benchmarks are relevant here?")

    return followups[:3]


# ═══════════════════════════════════════════════════════════════════════════════
# NOVA CHAT FEEDBACK STORE (in-memory + optional Supabase)
# ═══════════════════════════════════════════════════════════════════════════════

_chat_feedback: list[dict] = []
_chat_feedback_lock = threading.Lock()
_MAX_FEEDBACK_ENTRIES = 10000


def _store_feedback(
    message_id: str,
    rating: str,
    conversation_id: str = "",
    feedback_text: str = "",
) -> None:
    """Store chat message feedback.

    Args:
        message_id: Unique message identifier.
        rating: 'up' or 'down'.
        conversation_id: Optional conversation identifier.
        feedback_text: Optional text feedback from user.
    """
    entry = {
        "message_id": message_id,
        "rating": rating,
        "conversation_id": conversation_id,
        "feedback_text": feedback_text[:500] if feedback_text else "",
        "timestamp": datetime.datetime.utcnow().isoformat() + "Z",
    }
    with _chat_feedback_lock:
        _chat_feedback.append(entry)
        # Evict oldest entries if over limit
        if len(_chat_feedback) > _MAX_FEEDBACK_ENTRIES:
            _chat_feedback[:] = _chat_feedback[-_MAX_FEEDBACK_ENTRIES:]

    # Async persist to Supabase if available
    try:
        _persist_feedback_to_supabase(entry)
    except Exception as exc:
        logger.error("Failed to persist feedback to Supabase: %s", exc, exc_info=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SUPABASE CONVERSATION PERSISTENCE (write logic -- table must exist)
# ═══════════════════════════════════════════════════════════════════════════════

_SUPABASE_URL = os.environ.get("SUPABASE_URL") or ""
_SUPABASE_KEY = os.environ.get("SUPABASE_SERVICE_ROLE_KEY") or ""


def _supabase_rest(
    table: str,
    method: str = "GET",
    payload: Optional[Any] = None,
    params: str = "",
) -> Optional[Any]:
    """Make a REST call to the Supabase PostgREST API.

    Args:
        table: Table name.
        method: HTTP method (GET, POST, PATCH).
        payload: JSON body for POST/PATCH.
        params: Query string (e.g., '?conversation_id=eq.abc&order=timestamp.asc').

    Returns:
        Parsed JSON response or None on failure.
    """
    if not _SUPABASE_URL or not _SUPABASE_KEY:
        return None
    url = f"{_SUPABASE_URL.rstrip('/')}/rest/v1/{table}{params}"
    headers = {
        "apikey": _SUPABASE_KEY,
        "Authorization": f"Bearer {_SUPABASE_KEY}",
        "Content-Type": "application/json",
        "Prefer": "return=representation",
    }
    body = json.dumps(payload).encode() if payload else None
    req = urllib.request.Request(url, data=body, headers=headers, method=method)
    try:
        with urllib.request.urlopen(req, timeout=10) as resp:
            return json.loads(resp.read().decode())
    except (urllib.error.URLError, OSError, json.JSONDecodeError, ValueError) as exc:
        logger.error(
            "Supabase REST %s %s failed: %s", method, table, exc, exc_info=True
        )
        return None


def _save_conversation_turn(
    conversation_id: str,
    user_message: str,
    assistant_response: str,
    model_used: str = "",
    sources: Optional[list] = None,
    confidence: float = 0.0,
) -> None:
    """Persist a conversation turn using the document model (W-05 fix).

    Delegates to nova_persistence.append_message() which appends both
    user and assistant messages to the conversation's JSONB messages
    array.  Includes retry-with-backoff (C-08) and a background queue
    for failed writes.

    Args:
        conversation_id: Session/conversation UUID.
        user_message: The user's message.
        assistant_response: Nova's response text.
        model_used: LLM provider/model string.
        sources: List of data sources used.
        confidence: Response confidence score.
    """
    if not _SUPABASE_URL or not _SUPABASE_KEY:
        return

    try:
        from nova_persistence import append_message
    except ImportError:
        logger.error(
            "nova_persistence module not available; falling back to no-op",
            exc_info=True,
        )
        return

    # Append the user message
    if user_message:
        try:
            append_message(
                conversation_id,
                "user",
                user_message,
                user_id="anonymous",
            )
        except Exception as exc:
            logger.error(
                "Failed to persist user message for %s: %s",
                conversation_id,
                exc,
                exc_info=True,
            )

    # Append the assistant response
    if assistant_response:
        try:
            append_message(
                conversation_id,
                "assistant",
                assistant_response,
                user_id="anonymous",
                model_used=model_used,
                sources=sources,
                confidence=confidence,
            )
        except Exception as exc:
            logger.error(
                "Failed to persist assistant message for %s: %s",
                conversation_id,
                exc,
                exc_info=True,
            )


def _load_conversation_history(conversation_id: str) -> list[dict]:
    """Load conversation history from the JSONB messages array (W-05 fix).

    Returns messages in the format the widget expects, with backward
    compatibility for the legacy row-per-turn format.

    Args:
        conversation_id: Session/conversation UUID.

    Returns:
        List of message dicts with role, content, timestamp, etc.
    """
    try:
        from nova_persistence import load_conversation_messages
    except ImportError:
        logger.error(
            "nova_persistence module not available for loading history",
            exc_info=True,
        )
        return []

    try:
        messages = load_conversation_messages(conversation_id)
        if messages:
            return messages
    except Exception as exc:
        logger.error(
            "Failed to load conversation %s from document model: %s",
            conversation_id,
            exc,
            exc_info=True,
        )

    # Fallback: try the old REST row-per-turn query in case migration
    # has not run yet.
    try:
        params = (
            f"?conversation_id=eq.{urllib.parse.quote(conversation_id)}"
            f"&order=timestamp.asc"
            f"&limit=100"
        )
        result = _supabase_rest("nova_conversations", method="GET", params=params)
        if isinstance(result, list) and result:
            converted: list[dict] = []
            for row in result:
                user_msg = row.get("user_message") or ""
                asst_msg = row.get("assistant_response") or ""
                ts = row.get("timestamp") or ""
                if user_msg:
                    converted.append(
                        {"role": "user", "content": user_msg, "timestamp": ts}
                    )
                if asst_msg:
                    msg_obj: dict = {
                        "role": "assistant",
                        "content": asst_msg,
                        "timestamp": ts,
                    }
                    model = row.get("model_used") or ""
                    if model:
                        msg_obj["model_used"] = model
                    converted.append(msg_obj)
            return converted
    except Exception as fallback_exc:
        logger.error(
            "Legacy history fallback also failed for %s: %s",
            conversation_id,
            fallback_exc,
            exc_info=True,
        )

    return []


def _list_conversations(limit: int = 50) -> list[dict]:
    """List recent conversations from Supabase (W-05 fix).

    Uses the document model via nova_persistence to retrieve
    conversation summaries.

    Args:
        limit: Maximum number of conversations to return.

    Returns:
        List of conversation summary dicts.
    """
    try:
        from nova_persistence import list_conversations_summary

        return list_conversations_summary(limit=min(limit, 200))
    except ImportError:
        logger.error(
            "nova_persistence module not available for listing conversations",
            exc_info=True,
        )
    except Exception as exc:
        logger.error(
            "Failed to list conversations via document model: %s",
            exc,
            exc_info=True,
        )

    # Fallback to old REST query using document model columns
    params = (
        f"?select=id,title,updated_at,messages"
        f"&order=updated_at.desc"
        f"&limit={min(limit, 200)}"
    )
    result = _supabase_rest("nova_conversations", method="GET", params=params)
    if not isinstance(result, list):
        return []

    summaries: list[dict] = []
    for row in result:
        cid = row.get("id") or ""
        messages = row.get("messages") or []
        last_msg = ""
        if isinstance(messages, list):
            for msg in reversed(messages):
                if isinstance(msg, dict) and msg.get("role") == "user":
                    last_msg = (msg.get("content") or "")[:100]
                    break
        summaries.append(
            {
                "conversation_id": cid,
                "title": row.get("title") or "New Chat",
                "last_message": last_msg,
                "updated_at": row.get("updated_at") or "",
            }
        )
    return summaries[:limit]


def _persist_feedback_to_supabase(entry: dict) -> None:
    """Persist a feedback entry to the nova_feedback Supabase table.

    Args:
        entry: Feedback dict with message_id, rating, conversation_id, timestamp.
    """
    if not _SUPABASE_URL or not _SUPABASE_KEY:
        return
    _supabase_rest("nova_feedback", method="POST", payload=entry)


# ═══════════════════════════════════════════════════════════════════════════════
# ROI CALCULATOR BACKEND
# ═══════════════════════════════════════════════════════════════════════════════

# Conversion rates and time-to-fill by industry.  CPA and cost-per-hire are
# now sourced from benchmark_registry (single source of truth) via
# _get_roi_industry_benchmarks().
_ROI_CONV_RATE_TTF: dict[str, dict[str, float]] = {
    "technology": {"conv_rate": 0.05, "avg_ttf": 42},
    "healthcare": {"conv_rate": 0.045, "avg_ttf": 49},
    "finance": {"conv_rate": 0.04, "avg_ttf": 44},
    "retail": {"conv_rate": 0.08, "avg_ttf": 36},
    "manufacturing": {"conv_rate": 0.06, "avg_ttf": 40},
    "logistics": {"conv_rate": 0.07, "avg_ttf": 38},
    "hospitality": {"conv_rate": 0.09, "avg_ttf": 33},
    "construction": {"conv_rate": 0.055, "avg_ttf": 41},
    "overall": {"conv_rate": 0.05, "avg_ttf": 45},
}

# Maps frontend / legacy industry keys to benchmark_registry keys
_ROI_INDUSTRY_KEY_MAP: dict[str, str] = {
    "Technology": "technology",
    "Healthcare": "healthcare",
    "Finance": "finance",
    "Retail": "retail",
    "Manufacturing": "manufacturing",
    "Logistics": "logistics",
    "Hospitality": "hospitality",
    "Other": "overall",
    "tech_engineering": "technology",
    "healthcare_medical": "healthcare",
    "finance_banking": "finance",
    "retail_consumer": "retail",
    "construction_real_estate": "construction",
    "logistics_supply_chain": "logistics",
    "hospitality_travel": "hospitality",
    "general": "overall",
}


def _get_roi_industry_benchmarks(industry: str) -> dict[str, float]:
    """Get CPA, cost-per-hire, conversion rate, and TTF for an industry.

    Pulls CPA and cost-per-hire from benchmark_registry (single source of
    truth); conv_rate and avg_ttf come from the local _ROI_CONV_RATE_TTF table
    since the registry doesn't track those.

    Args:
        industry: Industry key from frontend or KB (e.g., "Technology",
                  "tech_engineering", "healthcare_medical").

    Returns:
        Dict with keys: cpa, conv_rate, avg_cph, avg_ttf.
    """
    registry_key = _ROI_INDUSTRY_KEY_MAP.get(industry, "overall")
    conv_ttf = _ROI_CONV_RATE_TTF.get(registry_key, _ROI_CONV_RATE_TTF["overall"])

    # CPA: use the industry-adjusted programmatic CPA from the registry
    bench = benchmark_registry.get_channel_benchmark("programmatic", registry_key)
    cpa = bench.get("cpa_adjusted") or bench.get("cpa") or 15.0

    # Cost-per-hire from registry
    avg_cph = benchmark_registry.get_cost_per_hire(registry_key)

    return {
        "cpa": cpa,
        "conv_rate": conv_ttf["conv_rate"],
        "avg_cph": avg_cph,
        "avg_ttf": conv_ttf["avg_ttf"],
    }


# ═══════════════════════════════════════════════════════════════════════════════
# COMPLIANCEGUARD: LLM-Powered Job Posting Compliance Analysis
# ═══════════════════════════════════════════════════════════════════════════════


def _analyze_compliance(data: dict) -> dict:
    """Analyze a job posting for compliance issues using LLM.

    Args:
        data: Dict with keys text (job description), job_title, location, industry.

    Returns:
        Dict with score, issues list, recommendations, and compliance_reference.
    """
    text = data.get("text") or ""
    job_title = data.get("job_title") or ""
    location = data.get("location") or ""
    industry = data.get("industry") or ""

    # Fetch EEOC/OFCCP compliance reference data
    eeoc_data: dict = {}
    try:
        from api_enrichment import fetch_eeoc_requirements

        eeoc_data = fetch_eeoc_requirements(industry=industry, location=location)
    except ImportError:
        logger.warning("api_enrichment not available for EEOC data")
    except Exception as exc:
        logger.error("EEOC requirements fetch failed: %s", exc, exc_info=True)

    if not text:
        no_text_result: dict = {
            "error": "Job description text is required",
            "score": 0,
            "issues": [],
            "recommendations": [],
        }
        if eeoc_data:
            no_text_result["compliance_reference"] = eeoc_data
        return no_text_result

    router = _lazy_llm_router()
    if not router:
        no_llm_result: dict = {
            "score": 0,
            "issues": [],
            "recommendations": ["LLM service unavailable. Please try again later."],
            "llm_available": False,
        }
        if eeoc_data:
            no_llm_result["compliance_reference"] = eeoc_data
        return no_llm_result

    task_type = getattr(router, "TASK_COMPLIANCE_CHECK", "compliance_check")
    prompt = (
        f"Analyze this job posting for compliance issues.\n\n"
        f"Job Title: {job_title}\nLocation: {location}\n\n"
        f"Job Description:\n{text[:3000]}\n\n"
        f"Check for: gender bias, age discrimination, disability bias, "
        f"pay transparency requirements, EEO/OFCCP compliance.\n\n"
        f"For each issue found, provide: the problematic phrase, why it's an issue, "
        f"suggested replacement, and applicable law/regulation.\n\n"
        f"Return valid JSON with this exact structure:\n"
        f'{{"score": <0-100 compliance score>, '
        f'"issues": [{{"phrase": "...", "reason": "...", "suggestion": "...", "law": "..."}}], '
        f'"recommendations": ["general recommendation 1", ...]}}'
    )

    try:
        # S50: 10s timeout + TASK_PLAN_STRUCTURED for JSON output (Gemini best
        # at structured compliance JSON). Falls back through chain.
        result = router.call_llm(
            messages=[{"role": "user", "content": prompt}],
            system_prompt=(
                "You are an employment law and recruiting compliance expert. "
                "Analyze job postings for legal compliance issues. "
                "Always return valid JSON. Be thorough but practical."
            ),
            task_type=getattr(router, "TASK_PLAN_STRUCTURED", task_type),
            max_tokens=1024,
            timeout_budget=10.0,
        )
        response_text = result.get("text") or ""

        # Try to parse JSON from the response
        try:
            # Strip markdown code fences if present
            cleaned = response_text.strip()
            if cleaned.startswith("```"):
                cleaned = re.sub(r"^```(?:json)?\s*", "", cleaned)
                cleaned = re.sub(r"\s*```$", "", cleaned)
            parsed = json.loads(cleaned)
            llm_result: dict = {
                "score": parsed.get("score") or 0,
                "issues": parsed.get("issues") or [],
                "recommendations": parsed.get("recommendations") or [],
                "llm_provider": result.get("provider") or "",
            }
            if eeoc_data:
                llm_result["compliance_reference"] = eeoc_data
            return llm_result
        except (json.JSONDecodeError, ValueError):
            # LLM returned non-JSON; wrap as narrative
            narrative_result: dict = {
                "score": 50,
                "issues": [],
                "recommendations": [response_text[:1000]],
                "llm_provider": result.get("provider") or "",
            }
            if eeoc_data:
                narrative_result["compliance_reference"] = eeoc_data
            return narrative_result
    except Exception as e:
        logger.error("Compliance LLM analysis failed: %s", e, exc_info=True)
        err_result: dict = {
            "score": 0,
            "issues": [],
            "recommendations": ["Analysis temporarily unavailable. Please try again."],
            "llm_available": False,
        }
        if eeoc_data:
            err_result["compliance_reference"] = eeoc_data
        return err_result


# ═══════════════════════════════════════════════════════════════════════════════
# COMPLIANCEGUARD AUDIT: Enhanced LLM-Powered Compliance Audit with Fix Mode
# ═══════════════════════════════════════════════════════════════════════════════


def _audit_complianceguard(data: dict) -> dict:
    """Run a comprehensive compliance audit on a job posting using LLM.

    Supports both analysis-only and fix_mode (rewrite the job description).

    Args:
        data: Dict with job_title, location, job_description, company_size,
              industry, and optional fix_mode boolean.

    Returns:
        Dict with score, risk_level, findings, recommendations,
        and optionally rewritten_description.
    """
    job_description: str = data.get("job_description") or ""
    job_title: str = data.get("job_title") or ""
    location: str = data.get("location") or ""
    company_size: str = data.get("company_size") or ""
    industry: str = data.get("industry") or ""
    fix_mode: bool = bool(data.get("fix_mode"))

    if not job_description:
        return {"error": "job_description is required", "score": 0, "findings": []}

    # Enrich with O*NET occupation classification
    onet_context: str = ""
    if _api_integrations_available and _api_onet and job_title:
        try:
            occ_results = _api_onet.search_occupations(job_title)
            if occ_results:
                top_occ = occ_results[0]
                soc_code = top_occ.get("code") or ""
                onet_context = (
                    f"O*NET Classification: {top_occ.get('title') or ''} "
                    f"(SOC: {soc_code})"
                )
                try:
                    skills = _api_onet.get_skills(soc_code)
                    if skills:
                        skill_names = [s.get("name") or "" for s in skills[:5]]
                        onet_context += f"\nKey skills: {', '.join(skill_names)}"
                except (TypeError, KeyError, ValueError) as sk_err:
                    logger.error(f"O*NET skills fetch failed: {sk_err}", exc_info=True)
        except (TypeError, KeyError, ValueError) as onet_err:
            logger.error(f"O*NET lookup failed: {onet_err}", exc_info=True)

    # Fetch Supabase compliance rules
    supabase_rules: list = []
    if _supabase_data_available and location:
        try:
            supabase_rules = get_compliance_rules(jurisdiction=location) or []
        except (urllib.error.URLError, OSError, ValueError) as sb_err:
            logger.error(
                f"Supabase compliance rules fetch failed: {sb_err}", exc_info=True
            )

    rules_context: str = ""
    if supabase_rules:
        rules_context = "\n\nApplicable Compliance Rules:\n"
        for rule in supabase_rules[:10]:
            rule_name = rule.get("name") or rule.get("rule_name") or ""
            rule_desc = rule.get("description") or rule.get("details") or ""
            if rule_name:
                rules_context += f"- {rule_name}: {rule_desc[:200]}\n"

    router = _lazy_llm_router()
    if not router:
        return {
            "score": 0,
            "risk_level": "unknown",
            "findings": [],
            "recommendations": ["LLM service unavailable. Please try again later."],
            "rewritten_description": None,
        }

    task_type = getattr(router, "TASK_COMPLIANCE_CHECK", "compliance_check")

    # Build the analysis + fix prompt
    fix_instruction: str = ""
    if fix_mode:
        fix_instruction = (
            '\n\nAlso provide a "rewritten_description" field with a fully '
            "compliant rewrite of the job description that fixes ALL identified "
            "issues. Keep the original meaning and requirements but make it "
            "legally compliant across all relevant jurisdictions."
        )

    prompt: str = (
        f"Perform a comprehensive compliance audit on this job posting.\n\n"
        f"Job Title: {job_title}\n"
        f"Location: {location}\n"
        f"Company Size: {company_size}\n"
        f"Industry: {industry}\n"
        f"{onet_context}\n{rules_context}\n\n"
        f"Job Description:\n{job_description[:4000]}\n\n"
        f"Analyze against ALL of the following compliance areas:\n"
        f"1. EEOC compliance (US Equal Employment Opportunity)\n"
        f"2. GDPR / data privacy (EU regulations, June 2026 deadline)\n"
        f"3. Pay transparency (NYC, CO, CA, WA, EU Pay Transparency Directive)\n"
        f"4. Age discrimination (ADEA protections)\n"
        f"5. Disability accommodation language (ADA)\n"
        f"6. Gender-neutral language check\n\n"
        f"For each finding, provide severity (Critical/High/Medium/Low), "
        f"the specific text flagged, the applicable law, and a recommended fix.\n\n"
        f"Return valid JSON with this exact structure:\n"
        f'{{"score": <0-100>, "risk_level": "<Low|Medium|High|Critical>", '
        f'"findings": [{{"category": "...", "severity": "...", "description": "...", '
        f'"flagged_text": "...", "applicable_law": "...", "recommended_fix": "..."}}], '
        f'"recommendations": ["general recommendation 1", ...]'
        f"{fix_instruction}"
        f"}}"
    )

    try:
        # S50: Route compliance audits to TASK_PLAN_STRUCTURED (Gemini -- best
        # structured JSON output). 10s timeout for non-fix mode; fix mode gets
        # 15s because it includes a full rewrite.
        _compliance_timeout = 15.0 if fix_mode else 10.0
        result = router.call_llm(
            messages=[{"role": "user", "content": prompt}],
            system_prompt=(
                "You are an expert employment law and recruiting compliance analyst "
                "specializing in US (EEOC, ADA, ADEA, OFCCP) and EU (GDPR, EU Pay "
                "Transparency Directive, EU Employment Equality Directive) regulations. "
                "Analyze job postings thoroughly. Always return valid JSON."
            ),
            task_type=getattr(router, "TASK_PLAN_STRUCTURED", task_type),
            max_tokens=2048 if fix_mode else 1536,
            timeout_budget=_compliance_timeout,
        )
        response_text: str = result.get("text") or ""

        # Parse JSON from response
        cleaned: str = response_text.strip()
        if cleaned.startswith("```"):
            cleaned = re.sub(r"^```(?:json)?\s*", "", cleaned)
            cleaned = re.sub(r"\s*```$", "", cleaned)

        try:
            parsed = json.loads(cleaned)
        except (json.JSONDecodeError, ValueError):
            # Fallback: wrap as narrative
            return {
                "score": 50,
                "risk_level": "Medium",
                "findings": [],
                "recommendations": [response_text[:1000]],
                "rewritten_description": None,
                "llm_provider": result.get("provider") or "",
            }

        audit_result: dict = {
            "score": parsed.get("score") or 0,
            "risk_level": parsed.get("risk_level") or "Medium",
            "findings": parsed.get("findings") or [],
            "recommendations": parsed.get("recommendations") or [],
            "rewritten_description": parsed.get("rewritten_description"),
            "llm_provider": result.get("provider") or "",
        }
        if onet_context:
            audit_result["onet_classification"] = onet_context
        if supabase_rules:
            audit_result["applicable_rules_count"] = len(supabase_rules)
        return audit_result

    except Exception as e:
        logger.error(f"ComplianceGuard audit LLM call failed: {e}", exc_info=True)
        return {
            "score": 0,
            "risk_level": "unknown",
            "findings": [],
            "recommendations": ["Analysis temporarily unavailable. Please try again."],
            "rewritten_description": None,
        }


# ═══════════════════════════════════════════════════════════════════════════════
# CREATIVEAI: LLM-Powered Ad Copy Generation
# ═══════════════════════════════════════════════════════════════════════════════


def _generate_creative_ads(data: dict) -> dict:
    """Generate recruitment ad copy variants using LLM.

    Args:
        data: Dict with keys job_title, company, selling_points, tone,
              platform, char_limit.

    Returns:
        Dict with variants list.
    """
    job_title = data.get("job_title") or ""
    company = data.get("company") or ""
    selling_points = data.get("selling_points") or ""
    tone = data.get("tone") or "professional"
    platform = data.get("platform") or "LinkedIn"
    char_limit = int(data.get("char_limit") or 150)

    if not job_title:
        return {"error": "job_title is required", "variants": []}

    # ── API enrichment: O*NET skills + Tavily trending topics ──
    skills_context = ""
    trending_context = ""
    if _api_integrations_available and _api_onet:
        try:
            occ_results = _api_onet.search_occupations(job_title)
            if occ_results:
                soc_code = occ_results[0].get("code") or ""
                if soc_code:
                    skills = _api_onet.get_skills(soc_code)
                    if skills:
                        top_skills = [
                            s.get("name") or "" for s in skills[:5] if s.get("name")
                        ]
                        if top_skills:
                            skills_context = (
                                f"Key skills for this role: {', '.join(top_skills)}\n"
                            )
        except Exception as api_err:
            logger.error(f"CreativeAI O*NET enrichment error: {api_err}", exc_info=True)
    if _tavily_available and _tavily_recruitment_news:
        try:
            trending = _tavily_recruitment_news(
                f"recruitment trends {job_title} hiring"
            )
            if trending and isinstance(trending, list):
                headlines = [
                    t.get("title") or "" for t in trending[:3] if t.get("title")
                ]
                if headlines:
                    trending_context = (
                        f"Current recruitment trends: {'; '.join(headlines)}\n"
                    )
        except Exception as api_err:
            logger.error(
                f"CreativeAI Tavily enrichment error: {api_err}", exc_info=True
            )

    router = _lazy_llm_router()
    if not router:
        return {"variants": [], "llm_available": False}

    # S48: Route ad copy (prose generation) to Groq via TASK_PLAN_NARRATIVE
    # S50: Output is JSON array -- use TASK_PLAN_STRUCTURED (Gemini) for reliable
    # structured output. Prose quality inside the JSON values is still good with
    # Gemini. 10s timeout to avoid blocking creative generation.
    task_type = getattr(router, "TASK_PLAN_STRUCTURED", "plan_structured")
    prompt = (
        f"Generate 5 recruitment ad copy variants for:\n"
        f"Role: {job_title} at {company}\n"
        f"Key selling points: {selling_points}\n"
    )
    if skills_context:
        prompt += skills_context
    if trending_context:
        prompt += trending_context
    prompt += (
        f"Tone: {tone}\n"
        f"Platform: {platform}\n"
        f"Max {char_limit} characters per field.\n\n"
        f"Create:\n"
        f"1) Benefits-focused\n"
        f"2) Culture-forward\n"
        f"3) Growth-focused\n"
        f"4) Impact-driven\n"
        f"5) Data-backed\n\n"
        f"Return valid JSON array of objects with keys: variant_name, headline, body, cta"
    )

    try:
        result = router.call_llm(
            messages=[{"role": "user", "content": prompt}],
            system_prompt=(
                "You are a top-tier recruitment copywriter. "
                "Write compelling, concise ad copy. Always return valid JSON array. "
                "Each variant must be distinct in approach."
            ),
            task_type=task_type,
            max_tokens=800,
            timeout_budget=10.0,
        )
        response_text = result.get("text") or ""

        try:
            cleaned = response_text.strip()
            if cleaned.startswith("```"):
                cleaned = re.sub(r"^```(?:json)?\s*", "", cleaned)
                cleaned = re.sub(r"\s*```$", "", cleaned)
            parsed = json.loads(cleaned)
            variants = (
                parsed if isinstance(parsed, list) else parsed.get("variants") or []
            )
            return {
                "variants": variants,
                "llm_provider": result.get("provider") or "",
            }
        except (json.JSONDecodeError, ValueError):
            return {
                "variants": [],
                "raw_response": response_text[:1500],
                "llm_provider": result.get("provider") or "",
            }
    except Exception as e:
        logger.error("Creative ad generation failed: %s", e, exc_info=True)
        return {"variants": [], "llm_available": False}


# ═══════════════════════════════════════════════════════════════════════════════
# POST-CAMPAIGN ANALYSIS: LLM-Powered Executive Summary
# ═══════════════════════════════════════════════════════════════════════════════


def _generate_post_campaign_summary(data: dict) -> dict:
    """Generate an AI-powered executive summary of campaign results.

    Args:
        data: Campaign metrics dict with channels, spend, hires, CPA, etc.

    Returns:
        Dict with executive_summary, channel_analysis, recommendations.
    """
    if not data:
        return {
            "error": "Campaign data is required",
            "executive_summary": "",
            "channel_analysis": "",
            "recommendations": "",
        }

    # ── API enrichment: FRED macro context + BLS projections ──
    macro_context_str = ""
    if _api_integrations_available:
        macro_parts: list[str] = []
        try:
            if _api_fred:
                unemp = _api_fred.get_unemployment_rate()
                if unemp:
                    rate = unemp.get("value") or unemp.get("rate") or ""
                    if rate:
                        macro_parts.append(f"Current US unemployment rate: {rate}%")
        except Exception as api_err:
            logger.error(
                f"Post-campaign FRED enrichment error: {api_err}", exc_info=True
            )
        try:
            job_title_raw = (
                data.get("job_title") or data.get("role") or data.get("roles") or ""
            )
            if isinstance(job_title_raw, list):
                job_title_raw = job_title_raw[0] if job_title_raw else ""
            if job_title_raw and _api_onet and _api_bls:
                occ_results = _api_onet.search_occupations(str(job_title_raw))
                if occ_results:
                    soc_code = occ_results[0].get("code") or ""
                    if soc_code:
                        soc_clean = soc_code.replace("-", "").replace(".", "")[:6]
                        projections = _api_bls.get_employment_projections(soc_clean)
                        if projections:
                            macro_parts.append(
                                f"BLS employment projections: {json.dumps(projections, default=str)[:200]}"
                            )
        except Exception as api_err:
            logger.error(
                f"Post-campaign BLS enrichment error: {api_err}", exc_info=True
            )
        if macro_parts:
            macro_context_str = "Market context: " + "; ".join(macro_parts) + "\n\n"

    router = _lazy_llm_router()
    if not router:
        return _post_campaign_template_fallback(data)

    # S48: Route executive summaries to Groq via TASK_PLAN_NARRATIVE (fast prose)
    # S50: Output is JSON with prose values -- use TASK_PLAN_STRUCTURED (Gemini)
    # for reliable JSON structure, 10s timeout.
    task_type = getattr(router, "TASK_PLAN_STRUCTURED", "plan_structured")
    prompt = f"Generate an executive summary of this recruitment campaign.\n\n"
    if macro_context_str:
        prompt += macro_context_str
    prompt += (
        f"Campaign Data:\n{json.dumps(data, indent=2, default=str)[:2500]}\n\n"
        f"Write as a senior recruitment strategist presenting to a VP of Talent Acquisition.\n\n"
        f"Analyze:\n"
        f"1) What worked well (which channels outperformed and WHY -- cite CPA, hires, conversion rates)\n"
        f"2) What underperformed (high CPA channels -- quantify the gap vs benchmark)\n"
        f"3) Specific recommendations for next campaign with projected impact\n"
        f"4) Key risks if recommendations are not followed\n\n"
        f"Every statement must include a specific number. No generic observations.\n\n"
        f"Return valid JSON with keys: executive_summary, channel_analysis, recommendations"
    )

    try:
        result = router.call_llm(
            messages=[{"role": "user", "content": prompt}],
            system_prompt=(
                "You are a senior recruitment marketing strategist presenting to C-suite executives. "
                "Write with authority and precision. Every sentence must cite a specific metric. "
                "Include market thesis, ROI impact, and risk analysis. Always return valid JSON."
            ),
            task_type=task_type,
            max_tokens=800,
            timeout_budget=10.0,
        )
        response_text = result.get("text") or ""

        try:
            cleaned = response_text.strip()
            if cleaned.startswith("```"):
                cleaned = re.sub(r"^```(?:json)?\s*", "", cleaned)
                cleaned = re.sub(r"\s*```$", "", cleaned)
            parsed = json.loads(cleaned)
            return {
                "executive_summary": parsed.get("executive_summary") or "",
                "channel_analysis": parsed.get("channel_analysis") or "",
                "recommendations": parsed.get("recommendations") or "",
                "llm_provider": result.get("provider") or "",
            }
        except (json.JSONDecodeError, ValueError):
            return {
                "executive_summary": response_text[:1000],
                "channel_analysis": "",
                "recommendations": "",
                "llm_provider": result.get("provider") or "",
            }
    except Exception as e:
        logger.error("Post-campaign summary generation failed: %s", e, exc_info=True)
        return _post_campaign_template_fallback(data)


def _post_campaign_template_fallback(data: dict) -> dict:
    """Build a template-based post-campaign summary when LLM is unavailable.

    Extracts key metrics from the campaign data dict and formats them into
    structured text, ensuring the UI always has content to display.
    """
    total_spend = 0.0
    total_hires = 0
    channels_data = data.get("channels") or data.get("channel_data") or []
    best_channel = ""
    best_cpa = float("inf")

    for ch in channels_data:
        if isinstance(ch, dict):
            spend = float(ch.get("spend") or ch.get("budget") or 0)
            hires = int(ch.get("hires") or ch.get("applications") or 0)
            total_spend += spend
            total_hires += hires
            cpa = spend / hires if hires > 0 else float("inf")
            name = ch.get("name") or ch.get("channel") or "Unknown"
            if cpa < best_cpa:
                best_cpa = cpa
                best_channel = name

    exec_parts: list[str] = []
    if total_spend > 0:
        exec_parts.append(f"Total campaign spend: ${total_spend:,.0f}")
    if total_hires > 0:
        exec_parts.append(f"Total hires/applications: {total_hires}")
        if total_spend > 0:
            exec_parts.append(f"Overall CPA: ${total_spend / total_hires:,.0f}")
    if best_channel and best_cpa < float("inf"):
        exec_parts.append(
            f"Best performing channel: {best_channel} (CPA: ${best_cpa:,.0f})"
        )

    exec_summary = (
        "AI narrative temporarily unavailable. Key metrics from campaign data:\n"
        + "\n".join(f"- {p}" for p in exec_parts)
        if exec_parts
        else "Executive summary unavailable -- see data sheets for detailed analysis."
    )

    return {
        "executive_summary": exec_summary,
        "channel_analysis": (
            "Channel-level analysis unavailable -- review individual channel "
            "metrics in the data export."
        ),
        "recommendations": (
            "AI-powered recommendations temporarily unavailable. Review channel "
            "CPA and conversion rates to identify optimization opportunities."
        ),
        "llm_available": False,
    }


# ═══════════════════════════════════════════════════════════════════════════════
# ROI CALCULATOR
# ═══════════════════════════════════════════════════════════════════════════════


def _calculate_roi(data: dict) -> dict:
    """Calculate ROI metrics from provided recruitment advertising parameters.

    Accepts a dict with keys: monthly_budget, avg_cpa, avg_cost_per_hire,
    current_time_to_fill, num_positions, industry. Returns calculated ROI
    metrics including savings, payback period, and industry comparison.
    """
    monthly_budget = float(data.get("monthly_budget") or 0)
    avg_cpa = float(data.get("avg_cpa") or 0)
    avg_cost_per_hire = float(data.get("avg_cost_per_hire") or 0)
    current_ttf = float(data.get("current_time_to_fill") or 45)
    num_positions = int(data.get("num_positions") or 1)
    industry = data.get("industry") or "Other"

    if monthly_budget <= 0:
        return {"error": "monthly_budget must be greater than zero"}

    # Look up industry benchmarks from benchmark_registry
    benchmarks = _get_roi_industry_benchmarks(industry)
    cpa = avg_cpa if avg_cpa > 0 else benchmarks["cpa"]
    conv_rate = benchmarks["conv_rate"]
    industry_avg_cph = benchmarks["avg_cph"]
    industry_avg_ttf = benchmarks["avg_ttf"]

    # Core calculations
    monthly_applications = monthly_budget / cpa if cpa > 0 else 0
    projected_hires = monthly_applications * conv_rate
    cost_per_hire = monthly_budget / projected_hires if projected_hires > 0 else 0

    # Time-to-fill improvement (programmatic typically cuts 25-35%)
    ttf_improvement_pct = 30.0
    new_ttf = current_ttf * (1 - ttf_improvement_pct / 100)

    # Annual savings: difference in cost-per-hire times positions times 12
    current_annual_hiring_cost = (
        avg_cost_per_hire * num_positions * 12
        if avg_cost_per_hire > 0
        else industry_avg_cph * num_positions * 12
    )
    projected_annual_hiring_cost = cost_per_hire * num_positions * 12
    annual_ad_spend = monthly_budget * 12
    annual_savings = (
        current_annual_hiring_cost - projected_annual_hiring_cost - annual_ad_spend
    )

    # ROI percentage
    roi_percentage = (
        (annual_savings / annual_ad_spend * 100) if annual_ad_spend > 0 else 0
    )

    # Payback period
    monthly_savings = annual_savings / 12 if annual_savings > 0 else 0
    payback_period_months = (
        round(annual_ad_spend / (monthly_savings * 12) * 12, 1)
        if monthly_savings > 0
        else 0
    )

    # Industry benchmark comparison
    cph_vs_benchmark = "below" if cost_per_hire < industry_avg_cph else "above"
    cph_diff_pct = (
        abs(cost_per_hire - industry_avg_cph) / industry_avg_cph * 100
        if industry_avg_cph > 0
        else 0
    )

    # Try to enrich with KB data
    kb_insights: list[str] = []
    try:
        kb = load_knowledge_base()
        kb_key = _INDUSTRY_KEY_TO_KB_KEY.get(industry) or ""
        ind_benchmarks = kb.get("industry_specific_benchmarks", {}).get(kb_key, {})
        if ind_benchmarks:
            outlook = (
                ind_benchmarks.get("outlook")
                or ind_benchmarks.get("hiring_stability")
                or ""
            )
            if outlook:
                kb_insights.append(outlook)
            hiring_strength = ind_benchmarks.get("hiring_strength") or ""
            if hiring_strength:
                kb_insights.append(hiring_strength)
    except Exception as e:
        logger.error("ROI KB enrichment error: %s", e, exc_info=True)

    # ── API enrichment: BLS wage data + FRED CPI ──
    api_wage_data: dict | None = None
    api_cpi_data: dict | None = None
    if _api_integrations_available:
        try:
            if _api_bls:
                # Search O*NET for SOC code, then fetch BLS wages
                job_title_raw = data.get("job_title") or data.get("role") or ""
                if job_title_raw and _api_onet:
                    occ_results = _api_onet.search_occupations(job_title_raw)
                    if occ_results:
                        soc_code = occ_results[0].get("code") or ""
                        if soc_code:
                            soc_clean = soc_code.replace("-", "").replace(".", "")[:6]
                            api_wage_data = _api_bls.get_occupational_employment(
                                soc_clean
                            )
        except Exception as e:
            logger.error(f"ROI BLS wage enrichment error: {e}", exc_info=True)
        try:
            if _api_fred:
                api_cpi_data = _api_fred.get_cpi_data(months=6)
        except Exception as e:
            logger.error(f"ROI FRED CPI enrichment error: {e}", exc_info=True)

    # Generate AI strategic recommendation
    strategic_recommendation = _generate_narrative(
        data={
            "roi_percentage": round(roi_percentage, 1),
            "payback_period_months": payback_period_months,
            "annual_savings": round(annual_savings, 2),
            "cost_per_hire": round(cost_per_hire, 2),
            "industry": industry,
            "num_positions": num_positions,
        },
        product_context=(
            f"ROI projections for {industry} recruitment. "
            f"ROI: {roi_percentage:.1f}%, payback: {payback_period_months} months, "
            f"annual savings: ${annual_savings:,.0f}. "
            f"Write a 2-sentence strategic recommendation for the hiring team."
        ),
        max_tokens=200,
    )

    return {
        "monthly_applications": round(monthly_applications, 1),
        "projected_hires": round(projected_hires, 1),
        "cost_per_hire": round(cost_per_hire, 2),
        "time_to_fill_improvement": round(ttf_improvement_pct, 1),
        "new_time_to_fill": round(new_ttf, 1),
        "annual_savings": round(annual_savings, 2),
        "monthly_savings": round(annual_savings / 12, 2) if annual_savings else 0,
        "roi_percentage": round(roi_percentage, 1),
        "payback_period_months": payback_period_months,
        "annual_ad_spend": round(annual_ad_spend, 2),
        "industry_benchmark_comparison": {
            "industry": industry,
            "industry_avg_cph": industry_avg_cph,
            "your_projected_cph": round(cost_per_hire, 2),
            "comparison": cph_vs_benchmark,
            "difference_pct": round(cph_diff_pct, 1),
            "industry_avg_ttf": industry_avg_ttf,
            "your_projected_ttf": round(new_ttf, 1),
        },
        "kb_insights": kb_insights,
        "strategic_recommendation": strategic_recommendation,
        "api_enrichment": {
            "bls_wage_data": api_wage_data,
            "cpi_inflation": api_cpi_data,
        },
        "inputs": {
            "monthly_budget": monthly_budget,
            "avg_cpa": cpa,
            "avg_cost_per_hire": avg_cost_per_hire or industry_avg_cph,
            "current_time_to_fill": current_ttf,
            "num_positions": num_positions,
            "industry": industry,
        },
    }


def _generate_roi_excel(data: dict) -> bytes:
    """Generate an Excel summary of ROI projection data.

    Accepts the same data dict as _calculate_roi and produces a formatted
    .xlsx workbook with inputs, projected metrics, and industry comparison.
    """
    # Calculate if not already calculated
    if "monthly_applications" not in data:
        data = _calculate_roi(data)

    wb = Workbook()
    ws = wb.active
    ws.title = "ROI Projection"

    header_font = Font(name="Calibri", bold=True, size=12, color="FFFFFF")
    header_fill = PatternFill(
        start_color="6366F1", end_color="6366F1", fill_type="solid"
    )
    subheader_font = Font(name="Calibri", bold=True, size=10)
    normal_font = Font(name="Calibri", size=10)

    # Title
    ws.merge_cells("A1:B1")
    ws["A1"] = "ROI Projection Calculator"
    ws["A1"].font = Font(name="Calibri", bold=True, size=16, color="6366F1")

    # Inputs section
    ws["A3"] = "INPUTS"
    ws["A3"].font = Font(name="Calibri", bold=True, size=12, color="6366F1")
    inputs = data.get("inputs", data)
    input_rows = [
        ("Monthly Budget", f"${inputs.get('monthly_budget') or 0:,.0f}"),
        ("Avg CPA", f"${inputs.get('avg_cpa') or 0:,.2f}"),
        ("Current Cost-Per-Hire", f"${inputs.get('avg_cost_per_hire') or 0:,.0f}"),
        (
            "Current Time-to-Fill",
            f"{inputs.get('current_time_to_fill') or 0:.0f} days",
        ),
        ("Number of Positions", str(inputs.get("num_positions") or 0)),
        ("Industry", str(inputs.get("industry") or "")),
    ]
    for i, (label, val) in enumerate(input_rows, start=4):
        ws.cell(row=i, column=1, value=label).font = subheader_font
        ws.cell(row=i, column=2, value=val).font = normal_font

    # Results section
    results_start = len(input_rows) + 5
    ws.cell(row=results_start, column=1, value="PROJECTED RESULTS").font = Font(
        name="Calibri", bold=True, size=12, color="6366F1"
    )
    result_rows = [
        ("Monthly Applications", f"{data.get('monthly_applications') or 0:,.1f}"),
        ("Projected Monthly Hires", f"{data.get('projected_hires') or 0:,.1f}"),
        ("Projected Cost-Per-Hire", f"${data.get('cost_per_hire') or 0:,.2f}"),
        (
            "Time-to-Fill Improvement",
            f"{data.get('time_to_fill_improvement') or 0:.1f}%",
        ),
        ("New Time-to-Fill", f"{data.get('new_time_to_fill') or 0:.0f} days"),
        ("Annual Ad Spend", f"${data.get('annual_ad_spend') or 0:,.0f}"),
        ("Annual Savings", f"${data.get('annual_savings') or 0:,.0f}"),
        ("ROI Percentage", f"{data.get('roi_percentage') or 0:.1f}%"),
        (
            "Payback Period",
            f"{data.get('payback_period_months') or 0:.1f} months",
        ),
    ]
    for i, (label, val) in enumerate(result_rows, start=results_start + 1):
        ws.cell(row=i, column=1, value=label).font = subheader_font
        ws.cell(row=i, column=2, value=val).font = normal_font

    # Industry comparison
    bench_start = results_start + len(result_rows) + 2
    ws.cell(row=bench_start, column=1, value="INDUSTRY COMPARISON").font = Font(
        name="Calibri", bold=True, size=12, color="6366F1"
    )
    bench = data.get("industry_benchmark_comparison", {})
    bench_rows = [
        ("Industry Avg CPH", f"${bench.get('industry_avg_cph') or 0:,.0f}"),
        ("Your Projected CPH", f"${bench.get('your_projected_cph') or 0:,.2f}"),
        (
            "Comparison",
            f"{bench.get('difference_pct') or 0:.1f}% {bench.get('comparison') or ''} benchmark",
        ),
        ("Industry Avg TTF", f"{bench.get('industry_avg_ttf') or 0} days"),
        ("Your Projected TTF", f"{bench.get('your_projected_ttf') or 0:.0f} days"),
    ]
    for i, (label, val) in enumerate(bench_rows, start=bench_start + 1):
        ws.cell(row=i, column=1, value=label).font = subheader_font
        ws.cell(row=i, column=2, value=val).font = normal_font

    ws.column_dimensions["A"].width = 28
    ws.column_dimensions["B"].width = 30

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ═══════════════════════════════════════════════════════════════════════════════
# A/B TEST LAB BACKEND
# ═══════════════════════════════════════════════════════════════════════════════

import random as _random

# Template-based ad copy data for A/B test generation
_AB_HEADLINES_BENEFITS = {
    "active": [
        "Take Your Career to the Next Level as a {job_title}",
        "Join Our Team: {job_title} -- Top Benefits Included",
        "Now Hiring: {job_title} with Industry-Leading Benefits",
    ],
    "passive": [
        "Discover Why Top {job_title}s Are Making the Switch",
        "What If Your {job_title} Role Came with Better Benefits?",
        "See What {job_title}s Get at {company}",
    ],
    "early_career": [
        "Start Your {job_title} Career with Unmatched Benefits",
        "Launch Your {job_title} Journey -- Benefits from Day 1",
        "{job_title} Opportunity: Build Your Future with Great Benefits",
    ],
}

_AB_HEADLINES_CULTURE = {
    "active": [
        "Join a Team Where {job_title}s Thrive",
        "{company} Is Hiring: A Culture Built for {job_title}s",
        "Work Where Innovation Meets Purpose -- {job_title}",
    ],
    "passive": [
        "Curious About a Better Work Culture? {job_title} at {company}",
        "What Makes {company} the Best Place for {job_title}s?",
        "Your Next {job_title} Role Could Be at a Company That Cares",
    ],
    "early_career": [
        "Start Your Career in a Team-First Culture -- {job_title}",
        "Join {company} as a {job_title}: Culture That Grows With You",
        "Great Teams Start Here -- {job_title} at {company}",
    ],
}

_AB_BENEFITS_BY_INDUSTRY = {
    "tech_engineering": [
        "flexible remote work",
        "equity & stock options",
        "unlimited PTO",
        "learning stipend",
        "top-tier health insurance",
        "401(k) matching",
    ],
    "healthcare_medical": [
        "sign-on bonus",
        "tuition reimbursement",
        "flexible scheduling",
        "comprehensive health coverage",
        "student loan assistance",
        "shift differentials",
    ],
    "finance_banking": [
        "competitive base salary",
        "performance bonuses",
        "hybrid work",
        "professional development",
        "comprehensive benefits",
        "wellness programs",
    ],
    "retail_consumer": [
        "employee discounts",
        "flexible scheduling",
        "career advancement",
        "health benefits",
        "paid training",
        "bonus programs",
    ],
    "general": [
        "competitive pay",
        "health & dental insurance",
        "paid time off",
        "career growth opportunities",
        "work-life balance",
        "team events",
    ],
}

_AB_CULTURE_PHRASES = {
    "tech_engineering": [
        "innovation-first environment",
        "collaborative engineering culture",
        "autonomy and ownership",
        "continuous learning mindset",
    ],
    "healthcare_medical": [
        "patient-centered mission",
        "supportive team environment",
        "making a real difference",
        "compassion-driven workplace",
    ],
    "finance_banking": [
        "integrity and excellence",
        "client-focused approach",
        "meritocratic growth",
        "analytical rigor",
    ],
    "retail_consumer": [
        "customer-first mindset",
        "energetic team atmosphere",
        "fast-paced growth",
        "community-driven values",
    ],
    "general": [
        "collaborative team culture",
        "innovation and growth",
        "inclusive workplace",
        "mission-driven work",
    ],
}

# Channel name -> benchmark_registry key mapping for A/B test CPA lookup.
# Uses benchmark_registry as single source of truth instead of hardcoded values.
_AB_CHANNEL_TO_REGISTRY: dict[str, str] = {
    "Indeed": "indeed",
    "LinkedIn": "linkedin",
    "Facebook": "meta_facebook",
    "Google": "google_ads",
    "ZipRecruiter": "ziprecruiter",
    "Programmatic": "programmatic",
}


def _get_ab_cpa_estimate(channel: str) -> float:
    """Get CPA estimate for A/B test calculations from benchmark_registry.

    Args:
        channel: Display name of the channel (e.g., "Indeed", "LinkedIn").

    Returns:
        CPA value as a float. Falls back to programmatic CPA if channel unknown.
    """
    registry_key = _AB_CHANNEL_TO_REGISTRY.get(channel, "programmatic")
    bench = benchmark_registry.get_channel_benchmark(registry_key)
    return bench.get("cpa") or 25.0


def _generate_ab_test(data: dict) -> dict:
    """Generate an A/B test configuration with ad copy variants.

    Accepts a dict with keys: job_title, company, industry, test_type.
    If ANTHROPIC_API_KEY is set, uses Claude API for sophisticated copy.
    Otherwise falls back to template-based generation.
    """
    job_title = data.get("job_title") or "Software Engineer"
    company = data.get("company") or "Your Company"
    industry = data.get("industry") or "general"
    test_type = data.get("test_type") or "benefits_vs_culture"
    channel = data.get("channel") or "Indeed"
    budget = float(data.get("budget") or 5000)
    audience = data.get("audience") or "active"

    # Normalize industry key
    industry_key = industry
    for k in _INDUSTRY_KEY_TO_KB_KEY:
        if k == industry:
            industry_key = k
            break

    # ── API enrichment: Adzuna job benchmarks + BLS employment context ──
    ab_market_context: dict = {}
    if _api_integrations_available:
        try:
            if _api_adzuna:
                job_results = _api_adzuna.search_jobs(
                    job_title, "us", results_per_page=3
                )
                if job_results:
                    ab_market_context["adzuna_job_sample"] = job_results[:3]
                salary_hist = _api_adzuna.get_salary_histogram(job_title)
                if salary_hist:
                    ab_market_context["salary_benchmark"] = salary_hist
        except Exception as api_err:
            logger.error(f"A/B Test Adzuna enrichment error: {api_err}", exc_info=True)
        try:
            if _api_onet:
                occ_results = _api_onet.search_occupations(job_title)
                if occ_results:
                    soc_code = occ_results[0].get("code") or ""
                    if soc_code and _api_bls:
                        soc_clean = soc_code.replace("-", "").replace(".", "")[:6]
                        emp_data = _api_bls.get_occupational_employment(soc_clean)
                        if emp_data:
                            ab_market_context["bls_employment"] = emp_data
        except Exception as api_err:
            logger.error(f"A/B Test BLS enrichment error: {api_err}", exc_info=True)

    # Check for Claude API
    anthropic_key = (os.environ.get("ANTHROPIC_API_KEY") or "").strip()
    if anthropic_key:
        try:
            result = _generate_ab_test_with_claude(
                job_title,
                company,
                industry_key,
                audience,
                channel,
                budget,
                anthropic_key,
            )
            if result:
                if ab_market_context:
                    result["market_context"] = ab_market_context
                return result
        except Exception as e:
            logger.error(
                "Claude A/B test generation failed, using templates: %s",
                e,
                exc_info=True,
            )

    # Template-based fallback
    result = _generate_ab_test_template(
        job_title, company, industry_key, audience, channel, budget
    )
    if ab_market_context:
        result["market_context"] = ab_market_context
    return result


def _generate_ab_test_with_claude(
    job_title: str,
    company: str,
    industry: str,
    audience: str,
    channel: str,
    budget: float,
    api_key: str,
) -> "dict | None":
    """Generate A/B test variants via LLM Router, with Anthropic API fallback.

    Tries the LLM Router first (free providers like Groq, Gemini, etc.),
    then falls back to the direct Anthropic API if the router is unavailable.
    Returns the result dict or None if all calls fail.
    """
    prompt = (
        f"Generate two recruitment ad variants for A/B testing:\n\n"
        f"Job Title: {job_title}\n"
        f"Company: {company}\n"
        f"Industry: {industry}\n"
        f"Target Audience: {audience} job seekers\n"
        f"Channel: {channel}\n"
        f"Monthly Budget: ${budget:,.0f}\n\n"
        f"Variant A should focus on BENEFITS (compensation, perks, work-life balance).\n"
        f"Variant B should focus on CULTURE (team, mission, growth, values).\n\n"
        f"For each variant provide:\n"
        f"1. A compelling headline (under 60 characters)\n"
        f"2. A description (2-3 sentences, under 150 characters)\n"
        f"3. A call-to-action (under 30 characters)\n\n"
        f"Respond in this exact JSON format:\n"
        f'{{"variant_a": {{"headline": "...", "description": "...", "cta": "..."}},\n'
        f' "variant_b": {{"headline": "...", "description": "...", "cta": "..."}}}}'
    )

    system_prompt = (
        "You are a recruitment advertising copywriter specializing in "
        "high-performance job ad copy. Generate compelling, specific ad "
        "variants that drive clicks and applications."
    )

    # ── Try LLM Router first (free providers) ──
    try:
        # S50: Route A/B copy to TASK_PLAN_STRUCTURED (Gemini -- best JSON output)
        # instead of TASK_NARRATIVE (Claude Haiku -- expensive for ad copy).
        # 10s timeout for non-blocking A/B generation.
        from llm_router import call_llm, TASK_PLAN_STRUCTURED

        messages = [{"role": "user", "content": prompt}]
        llm_result = call_llm(
            messages=messages,
            system_prompt=system_prompt,
            task_type=TASK_PLAN_STRUCTURED,
            max_tokens=500,
            timeout_budget=10.0,
        )
        content_text = (
            llm_result.get("text") or ""
            if isinstance(llm_result, dict)
            else str(llm_result or "")
        )
        if content_text:
            json_match = re.search(r"\{[\s\S]*\}", content_text)
            if json_match:
                variants = json.loads(json_match.group())
                return _build_ab_response(
                    variants.get("variant_a", {}),
                    variants.get("variant_b", {}),
                    channel,
                    budget,
                    "llm_router",
                )
    except ImportError:
        logger.warning("llm_router module not available, falling back to direct API")
    except (json.JSONDecodeError, ValueError, TypeError) as e:
        logger.error("LLM Router A/B test parse failed: %s", e, exc_info=True)
    except Exception as e:
        logger.error("LLM Router A/B test failed: %s", e, exc_info=True)

    # ── Fallback: direct Anthropic API ──
    try:
        req_body = json.dumps(
            {
                "model": "claude-sonnet-4-20250514",
                "max_tokens": 500,
                "messages": [{"role": "user", "content": prompt}],
            }
        ).encode("utf-8")

        req = urllib.request.Request(
            "https://api.anthropic.com/v1/messages",
            data=req_body,
            headers={
                "Content-Type": "application/json",
                "x-api-key": api_key,
                "anthropic-version": "2023-06-01",
            },
            method="POST",
        )
        with urllib.request.urlopen(req, timeout=15) as resp:
            resp_data = json.loads(resp.read().decode("utf-8"))

        content_text = ""
        for block in resp_data.get("content") or []:
            if block.get("type") == "text":
                content_text += block.get("text") or ""

        json_match = re.search(r"\{[\s\S]*\}", content_text)
        if json_match:
            variants = json.loads(json_match.group())
            return _build_ab_response(
                variants.get("variant_a", {}),
                variants.get("variant_b", {}),
                channel,
                budget,
                "claude",
            )
    except (urllib.error.URLError, json.JSONDecodeError, KeyError) as e:
        logger.error("Claude API call failed: %s", e, exc_info=True)
    return None


def _generate_ab_test_template(
    job_title: str,
    company: str,
    industry: str,
    audience: str,
    channel: str,
    budget: float,
) -> dict:
    """Generate A/B test variants using templates (no API required).

    Returns a complete A/B test configuration dict with variants,
    statistical requirements, and recommendations.
    """
    ben_headlines = _AB_HEADLINES_BENEFITS.get(
        audience, _AB_HEADLINES_BENEFITS["active"]
    )
    cul_headlines = _AB_HEADLINES_CULTURE.get(audience, _AB_HEADLINES_CULTURE["active"])

    benefits = _AB_BENEFITS_BY_INDUSTRY.get(
        industry, _AB_BENEFITS_BY_INDUSTRY["general"]
    )
    cultures = _AB_CULTURE_PHRASES.get(industry, _AB_CULTURE_PHRASES["general"])

    headline_a = (
        _random.choice(ben_headlines)
        .replace("{job_title}", job_title)
        .replace("{company}", company)
    )
    headline_b = (
        _random.choice(cul_headlines)
        .replace("{job_title}", job_title)
        .replace("{company}", company)
    )

    selected_benefits = _random.sample(benefits, min(3, len(benefits)))
    selected_cultures = _random.sample(cultures, min(2, len(cultures)))

    desc_a = (
        f"Join {company} as a {job_title}. Enjoy {selected_benefits[0]}, "
        f"{selected_benefits[1]}, and {selected_benefits[2]}."
    )
    desc_b = (
        f"Be part of {company}'s {selected_cultures[0]}. As a {job_title}, "
        f"you will experience {selected_cultures[1]}."
    )

    cta_a = _random.choice(
        [
            "Apply Now -- See Full Benefits",
            "View Benefits & Apply",
            "Explore This Opportunity",
        ]
    )
    cta_b = _random.choice(
        [
            "Meet the Team & Apply",
            "Discover Our Culture",
            "See Why We Are Different",
        ]
    )

    variant_a = {"headline": headline_a, "description": desc_a, "cta": cta_a}
    variant_b = {"headline": headline_b, "description": desc_b, "cta": cta_b}

    return _build_ab_response(variant_a, variant_b, channel, budget, "template")


def _build_ab_response(
    variant_a: dict,
    variant_b: dict,
    channel: str,
    budget: float,
    source: str,
) -> dict:
    """Build the full A/B test response including statistical requirements.

    Combines variant ad copy with sample size calculations, recommended
    duration, and success metrics into a single response dict.
    """
    est_cpa = _get_ab_cpa_estimate(channel)
    daily_apps = max(1, round(budget / 30 / est_cpa))

    baseline_apply_rate = 0.05
    min_detectable_effect = 0.02
    min_sample = int(
        16
        * baseline_apply_rate
        * (1 - baseline_apply_rate)
        / (min_detectable_effect**2)
    )
    days_needed = max(7, min(int(min_sample / daily_apps) + 1, 90))

    split = "50/50" if budget >= 3000 else "60/40"

    return {
        "variant_a": variant_a,
        "variant_b": variant_b,
        "statistical_requirements": {
            "sample_size": min_sample,
            "confidence_level": "95%",
            "min_detectable_effect": f"{min_detectable_effect * 100:.0f}%",
            "baseline_apply_rate": f"{baseline_apply_rate * 100:.1f}%",
        },
        "recommended_duration_days": days_needed,
        "traffic_split": split,
        "estimated_daily_applications": daily_apps,
        "estimated_cpa": est_cpa,
        "success_metrics": [
            {"name": "CTR", "description": "Click-Through Rate"},
            {"name": "Apply Rate", "description": "Applications per Click"},
            {"name": "CPA", "description": "Cost per Application"},
            {"name": "Quality Score", "description": "Qualified Applicant Ratio"},
        ],
        "source": source,
        "channel": channel,
    }


# ═══════════════════════════════════════════════════════════════════════════════
# POST-CAMPAIGN PPT GENERATOR
# ═══════════════════════════════════════════════════════════════════════════════


def _generate_post_campaign_ppt(data: dict) -> bytes:
    """Generate a PowerPoint presentation for post-campaign analysis.

    Accepts a dict with campaign_name, client_name, start_date, end_date,
    total_budget, industry, benchmark_cpa, channels (list), overall_grade,
    and recommendations (list). Returns raw .pptx bytes.
    """
    from pptx import Presentation as PptxPresentation
    from pptx.util import Inches as PptxInches, Pt as PptxPt
    from pptx.dml.color import RGBColor as PptxRGB
    from pptx.enum.text import PP_ALIGN as PptxAlign

    prs = PptxPresentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    # Brand colors
    navy = PptxRGB(0x20, 0x20, 0x58)
    accent = PptxRGB(0x5A, 0x54, 0xBD)
    white = PptxRGB(0xFF, 0xFF, 0xFF)
    dark_text = PptxRGB(0x20, 0x20, 0x58)
    muted_text = PptxRGB(0x59, 0x67, 0x80)
    teal = PptxRGB(0x6B, 0xB3, 0xCD)

    campaign_name = data.get("campaign_name") or "Campaign Analysis"
    client_name = data.get("client_name") or ""
    start_date = data.get("start_date") or ""
    end_date = data.get("end_date") or ""
    benchmark_cpa = float(data.get("benchmark_cpa") or 20)
    channels = data.get("channels") or []
    overall_grade = data.get("overall_grade") or "--"
    recommendations = data.get("recommendations") or []

    def _add_bg(slide, color):
        """Set slide background color."""
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    # ── Slide 1: Title ──
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    _add_bg(slide1, navy)

    txBox = slide1.shapes.add_textbox(
        PptxInches(1), PptxInches(2), PptxInches(11), PptxInches(1.5)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Post-Campaign Analysis"
    p.font.size = PptxPt(40)
    p.font.bold = True
    p.font.color.rgb = white
    p.alignment = PptxAlign.LEFT

    p2 = tf.add_paragraph()
    p2.text = campaign_name
    p2.font.size = PptxPt(28)
    p2.font.color.rgb = teal
    p2.alignment = PptxAlign.LEFT

    if client_name or start_date:
        p3 = tf.add_paragraph()
        date_range = f"{start_date} -- {end_date}" if start_date and end_date else ""
        p3.text = f"{client_name}  |  {date_range}".strip(" | ")
        p3.font.size = PptxPt(16)
        p3.font.color.rgb = PptxRGB(0x8C, 0x96, 0xA8)
        p3.alignment = PptxAlign.LEFT

    brand_box = slide1.shapes.add_textbox(
        PptxInches(1), PptxInches(6.5), PptxInches(5), PptxInches(0.5)
    )
    bp = brand_box.text_frame.paragraphs[0]
    bp.text = "Nova AI Suite | Post-Campaign Analysis"
    bp.font.size = PptxPt(10)
    bp.font.color.rgb = PptxRGB(0x8C, 0x96, 0xA8)

    # ── Slide 2: Performance Summary ──
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide2, PptxRGB(0xFF, 0xFD, 0xF9))

    title2 = slide2.shapes.add_textbox(
        PptxInches(0.8), PptxInches(0.4), PptxInches(11), PptxInches(0.8)
    )
    t2p = title2.text_frame.paragraphs[0]
    t2p.text = "Performance Summary"
    t2p.font.size = PptxPt(28)
    t2p.font.bold = True
    t2p.font.color.rgb = dark_text

    total_spend = sum(float(c.get("spend") or 0) for c in channels)
    total_apps = sum(int(c.get("applications") or 0) for c in channels)
    total_hires = sum(int(c.get("hires") or 0) for c in channels)
    total_clicks = sum(int(c.get("clicks") or 0) for c in channels)
    total_impressions = sum(int(c.get("impressions") or 0) for c in channels)
    overall_cpa = total_spend / total_apps if total_apps > 0 else 0
    overall_cph = total_spend / total_hires if total_hires > 0 else 0
    overall_ctr = (
        (total_clicks / total_impressions * 100) if total_impressions > 0 else 0
    )

    metrics = [
        ("Overall Grade", overall_grade),
        ("Total Spend", f"${total_spend:,.0f}"),
        ("Total Applications", f"{total_apps:,}"),
        ("Total Hires", f"{total_hires:,}"),
        ("Overall CPA", f"${overall_cpa:,.2f}"),
        ("Overall CPH", f"${overall_cph:,.2f}"),
        ("Benchmark CPA", f"${benchmark_cpa:,.0f}"),
        ("CTR", f"{overall_ctr:.2f}%"),
    ]

    for i, (label, value) in enumerate(metrics):
        col = i % 4
        row = i // 4
        x = PptxInches(0.8 + col * 3.0)
        y = PptxInches(1.8 + row * 2.2)

        box = slide2.shapes.add_textbox(x, y, PptxInches(2.6), PptxInches(1.8))
        btf = box.text_frame
        btf.word_wrap = True

        lp = btf.paragraphs[0]
        lp.text = label
        lp.font.size = PptxPt(12)
        lp.font.color.rgb = muted_text

        vp = btf.add_paragraph()
        vp.text = str(value)
        vp.font.size = PptxPt(24)
        vp.font.bold = True
        vp.font.color.rgb = dark_text

    # ── Slide 3: Channel Breakdown ──
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide3, PptxRGB(0xFF, 0xFD, 0xF9))

    title3 = slide3.shapes.add_textbox(
        PptxInches(0.8), PptxInches(0.4), PptxInches(11), PptxInches(0.8)
    )
    t3p = title3.text_frame.paragraphs[0]
    t3p.text = "Channel Performance Breakdown"
    t3p.font.size = PptxPt(28)
    t3p.font.bold = True
    t3p.font.color.rgb = dark_text

    if channels:
        headers = [
            "Channel",
            "Spend",
            "Applications",
            "Hires",
            "CPA",
            "CPH",
            "Grade",
        ]
        rows_count = len(channels) + 1
        cols_count = len(headers)

        table_shape = slide3.shapes.add_table(
            rows_count,
            cols_count,
            PptxInches(0.5),
            PptxInches(1.5),
            PptxInches(12),
            PptxInches(0.4 * rows_count + 0.2),
        )
        table = table_shape.table

        for ci, hdr in enumerate(headers):
            cell = table.cell(0, ci)
            cell.text = hdr
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = PptxPt(11)
                paragraph.font.bold = True
                paragraph.font.color.rgb = white
            cell.fill.solid()
            cell.fill.fore_color.rgb = accent

        for ri, ch in enumerate(channels, start=1):
            ch_spend = float(ch.get("spend") or 0)
            ch_apps = int(ch.get("applications") or 0)
            ch_hires = int(ch.get("hires") or 0)
            ch_cpa = ch_spend / ch_apps if ch_apps > 0 else 0
            ch_cph = ch_spend / ch_hires if ch_hires > 0 else 0
            ch_grade = ch.get("grade") or "--"

            vals = [
                ch.get("channel") or "",
                f"${ch_spend:,.0f}",
                f"{ch_apps:,}",
                f"{ch_hires:,}",
                f"${ch_cpa:,.2f}",
                f"${ch_cph:,.2f}",
                ch_grade,
            ]
            for ci, val in enumerate(vals):
                cell = table.cell(ri, ci)
                cell.text = str(val)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = PptxPt(10)
                    paragraph.font.color.rgb = dark_text

    # ── Slide 4: Recommendations ──
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide4, navy)

    title4 = slide4.shapes.add_textbox(
        PptxInches(0.8), PptxInches(0.4), PptxInches(11), PptxInches(0.8)
    )
    t4p = title4.text_frame.paragraphs[0]
    t4p.text = "Recommendations"
    t4p.font.size = PptxPt(28)
    t4p.font.bold = True
    t4p.font.color.rgb = white

    rec_box = slide4.shapes.add_textbox(
        PptxInches(1), PptxInches(1.5), PptxInches(11), PptxInches(5)
    )
    rec_tf = rec_box.text_frame
    rec_tf.word_wrap = True

    if recommendations:
        for i, rec in enumerate(recommendations):
            rp = rec_tf.paragraphs[0] if i == 0 else rec_tf.add_paragraph()
            rec_text = rec if isinstance(rec, str) else str(rec)
            rp.text = f"  {rec_text}"
            rp.font.size = PptxPt(14)
            rp.font.color.rgb = PptxRGB(0xD4, 0xD4, 0xD8)
            rp.space_after = PptxPt(12)
    else:
        auto_recs = []
        if overall_cpa > benchmark_cpa * 1.3:
            auto_recs.append(
                f"Overall CPA (${overall_cpa:.2f}) exceeds benchmark "
                f"(${benchmark_cpa}) by "
                f"{((overall_cpa / benchmark_cpa) - 1) * 100:.0f}%. "
                f"Consider optimizing targeting and ad copy."
            )
        elif overall_cpa <= benchmark_cpa * 0.7:
            auto_recs.append(
                f"CPA (${overall_cpa:.2f}) is well below benchmark. "
                f"Consider scaling budget to capture more volume."
            )
        if channels:
            best = min(
                channels,
                key=lambda c: (
                    float(c.get("cpa") or 9999)
                    if float(c.get("cpa") or 0) > 0
                    else 9999
                ),
            )
            worst = max(channels, key=lambda c: float(c.get("cpa") or 0))
            if best.get("channel") != worst.get("channel"):
                auto_recs.append(
                    f"Shift budget from "
                    f"{worst.get('channel', 'underperforming channel')} toward "
                    f"{best.get('channel', 'top channel')} which has the "
                    f"lowest CPA."
                )
        auto_recs.append(
            "Implement A/B testing on ad creative to continuously "
            "improve conversion rates."
        )
        auto_recs.append(
            "Track quality-of-hire metrics alongside cost metrics "
            "for a complete performance picture."
        )

        for i, rec in enumerate(auto_recs):
            rp = rec_tf.paragraphs[0] if i == 0 else rec_tf.add_paragraph()
            rp.text = f"  {rec}"
            rp.font.size = PptxPt(14)
            rp.font.color.rgb = PptxRGB(0xD4, 0xD4, 0xD8)
            rp.space_after = PptxPt(12)

    brand4 = slide4.shapes.add_textbox(
        PptxInches(1), PptxInches(6.5), PptxInches(5), PptxInches(0.5)
    )
    b4p = brand4.text_frame.paragraphs[0]
    b4p.text = "Nova AI Suite | Post-Campaign Analysis"
    b4p.font.size = PptxPt(10)
    b4p.font.color.rgb = PptxRGB(0x8C, 0x96, 0xA8)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ═══════════════════════════════════════════════════════════════════════════════
# API HARDENING: Consistent Error Envelope
# ═══════════════════════════════════════════════════════════════════════════════
def _error_response(
    message: str, code: str = "UNKNOWN_ERROR", status: int = 400
) -> tuple[bytes, int]:
    """Build a consistent JSON error response.

    Returns a tuple of (encoded_body, http_status_code) for use in endpoint handlers.

    Args:
        message: Human-readable error description.
        code: Machine-readable error code (e.g., VALIDATION_ERROR, RATE_LIMITED).
        status: HTTP status code.
    """
    body = json.dumps(
        {
            "success": False,
            "error": message,
            "code": code,
            "data": None,
            "timestamp": datetime.datetime.now(datetime.timezone.utc).isoformat(),
        }
    )
    return body.encode("utf-8"), status


# ═══════════════════════════════════════════════════════════════════════════════
# API HARDENING: Request Deduplication (Issue 3)
# ═══════════════════════════════════════════════════════════════════════════════
_recent_requests: dict[str, float] = {}  # hash -> timestamp
_dedup_lock = threading.Lock()
_DEDUP_WINDOW = 60  # seconds


def _is_duplicate_request(payload: dict) -> bool:
    """Check if an identical request was made in the last 60 seconds.

    Hashes the payload (excluding timestamps/session IDs) and checks against
    a time-windowed cache. Thread-safe.

    Args:
        payload: The request JSON payload to check.

    Returns:
        True if a duplicate was found within the dedup window.
    """
    dedup_payload = {
        k: v
        for k, v in payload.items()
        if k not in ("_request_id", "session_id", "timestamp", "conversation_id")
    }
    dedup_key = hashlib.sha256(
        json.dumps(dedup_payload, sort_keys=True, default=str).encode()
    ).hexdigest()[:16]
    now = time.time()
    with _dedup_lock:
        # Clean old entries
        expired = [k for k, t in _recent_requests.items() if now - t > _DEDUP_WINDOW]
        for k in expired:
            del _recent_requests[k]
        # Check for duplicate
        if dedup_key in _recent_requests:
            return True
        _recent_requests[dedup_key] = now
        return False


# ═══════════════════════════════════════════════════════════════════════════════
# API HARDENING: Idempotency Key Support (P3 Improvement)
# ═══════════════════════════════════════════════════════════════════════════════
# Stores {idempotency_key: {response_body: bytes, status_code: int, expires: float}}
_idempotency_cache: dict[str, dict] = {}
_idempotency_lock = threading.Lock()
_IDEMPOTENCY_TTL_SECONDS: float = 300.0  # 5 minutes

# Endpoints that support idempotency keys
_IDEMPOTENT_ENDPOINTS: frozenset[str] = frozenset(
    {
        "/api/chat",
        "/api/chat/stream",
        "/api/generate",
        "/api/plan/share",
    }
)


def _idempotency_get(key: str) -> Optional[dict]:
    """Retrieve a cached response for an idempotency key if not expired.

    Thread-safe. Evicts expired entries opportunistically.

    Args:
        key: The idempotency key from X-Idempotency-Key header.

    Returns:
        Cached response dict with response_body and status_code, or None.
    """
    now = time.time()
    with _idempotency_lock:
        # Opportunistic cleanup of expired entries (cap at 50 to avoid latency)
        expired = [
            k
            for k, v in list(_idempotency_cache.items())[:200]
            if now > v.get("expires", 0)
        ]
        for k in expired[:50]:
            _idempotency_cache.pop(k, None)
        entry = _idempotency_cache.get(key)
        if entry and now <= entry.get("expires", 0):
            return entry
        # Expired or missing
        if entry:
            _idempotency_cache.pop(key, None)
        return None


def _idempotency_store(key: str, response_body: bytes, status_code: int) -> None:
    """Cache a response for an idempotency key with TTL.

    Thread-safe. Caps total cache size at 10,000 entries.

    Args:
        key: The idempotency key from X-Idempotency-Key header.
        response_body: The serialized response bytes.
        status_code: The HTTP status code of the response.
    """
    now = time.time()
    with _idempotency_lock:
        # Prevent unbounded growth
        if len(_idempotency_cache) >= 10_000:
            # Evict oldest 1000 entries
            sorted_keys = sorted(
                _idempotency_cache.keys(),
                key=lambda k: _idempotency_cache[k].get("expires", 0),
            )
            for k in sorted_keys[:1000]:
                _idempotency_cache.pop(k, None)
        _idempotency_cache[key] = {
            "response_body": response_body,
            "status_code": status_code,
            "expires": now + _IDEMPOTENCY_TTL_SECONDS,
        }


# ═══════════════════════════════════════════════════════════════════════════════
# API HARDENING: Generation Timeout (Issue 4)
# ═══════════════════════════════════════════════════════════════════════════════
_GENERATE_TIMEOUT_SECONDS = 240  # 4 minutes hard timeout


class _GenerationTimeoutError(Exception):
    """Raised when /api/generate exceeds the hard timeout."""

    pass


class MediaPlanHandler(BaseHTTPRequestHandler):
    def log_message(self, format, *args):
        logger.info(format, *args)

    def end_headers(self):
        """Add security + API platform headers to all responses."""
        # ── API Platform headers ──
        self.send_header("X-API-Version", "v1")
        # Request ID: generated per-request in do_GET/do_POST, or fallback
        _rid = getattr(self, "_request_id", None)
        if _rid:
            self.send_header("X-Request-ID", _rid)
        self.send_header("X-Content-Type-Options", "nosniff")
        self.send_header("X-Frame-Options", "SAMEORIGIN")
        self.send_header("Referrer-Policy", "strict-origin-when-cross-origin")
        # Prevent browsers from running inline scripts from injected content
        self.send_header("X-XSS-Protection", "1; mode=block")
        # Strict-Transport-Security: tell browsers to only use HTTPS
        self.send_header(
            "Strict-Transport-Security", "max-age=31536000; includeSubDomains"
        )
        # Permissions-Policy: disable sensitive browser APIs not needed by this app
        self.send_header(
            "Permissions-Policy", "camera=(), microphone=(), geolocation=()"
        )
        # Content-Security-Policy: restrict resource origins (unsafe-inline needed for
        # inline styles/scripts used throughout templates)
        self.send_header(
            "Content-Security-Policy",
            "default-src 'self'; "
            "script-src 'self' 'unsafe-inline' https://cdnjs.cloudflare.com https://unpkg.com https://cdn.jsdelivr.net https://us-assets.i.posthog.com https://us.i.posthog.com; "
            "style-src 'self' 'unsafe-inline' https://fonts.googleapis.com; "
            "font-src https://fonts.gstatic.com; "
            "img-src 'self' data: blob:; "
            "connect-src 'self' https://*.supabase.co https://*.sentry.io https://*.posthog.com https://us.i.posthog.com; "
            "object-src 'none'; "
            "base-uri 'self'; "
            "form-action 'self';",
        )
        super().end_headers()

    # ── PostHog Analytics Helpers ──────────────────────────────────────────

    def _ph_distinct_id(self) -> str:
        """Return a hashed distinct_id for the current request's client IP."""
        if not _posthog_available:
            return ""
        return _ph_hash_ip(self._get_client_ip())

    def _ph_track_page(self, path: str, template: str) -> None:
        """Track a page view event in PostHog (fire-and-forget)."""
        if not _posthog_available:
            return
        try:
            _ph_page_view(
                self._ph_distinct_id(),
                path,
                {
                    "template": template,
                    "user_agent": self.headers.get("User-Agent") or "",
                    "referer": self.headers.get("Referer") or "",
                },
            )
        except Exception as exc:
            logger.error("PostHog page_view tracking failed: %s", exc, exc_info=True)

    def _ph_track(self, event: str, properties: Optional[dict] = None) -> None:
        """Track a named event in PostHog (fire-and-forget)."""
        if not _posthog_available:
            return
        try:
            _ph_track(self._ph_distinct_id(), event, properties)
        except Exception as exc:
            logger.error("PostHog event tracking failed: %s", exc, exc_info=True)

    def _get_client_ip(self) -> str:
        """Return the real client IP, respecting reverse-proxy headers.

        On Render.com (behind a load balancer), ``self.client_address[0]``
        may return the proxy IP.  The ``X-Forwarded-For`` header carries the
        chain of client IPs.  The rightmost entry is the one appended by the
        trusted load balancer, so we use that to prevent spoofing via a
        client-supplied fake X-Forwarded-For header.
        """
        forwarded: str = self.headers.get("X-Forwarded-For") or ""
        if forwarded:
            # Use rightmost IP (set by trusted Render LB) to prevent spoofing
            return forwarded.split(",")[-1].strip()
        return self.client_address[0]

    def _get_cors_origin(self):
        """Return the request Origin if it is in the allowlist, else empty string.
        SECURITY: CORS_ALLOW_ALL removed to prevent accidental exposure in production.
        """
        origin = self.headers.get("Origin") or ""
        if origin in _ALLOWED_ORIGINS:
            return origin
        return ""  # No CORS header = browser blocks

    def _send_error(
        self, message: str, code: str = "UNKNOWN_ERROR", status: int = 400
    ) -> None:
        """Send a consistent JSON error response with CORS headers.

        Uses the module-level _error_response() to build the envelope, then
        writes it out with correct headers.

        Args:
            message: Human-readable error description.
            code: Machine-readable error code.
            status: HTTP status code.
        """
        body, status_code = _error_response(message, code, status)
        self.send_response(status_code)
        self.send_header("Content-Type", "application/json")
        cors_origin = self._get_cors_origin()
        if cors_origin:
            self.send_header("Access-Control-Allow-Origin", cors_origin)
        self.send_header("Content-Length", str(len(body)))
        self.end_headers()
        self.wfile.write(body)

    def _check_idempotency(self, path: str) -> bool:
        """Check for X-Idempotency-Key header and return cached response if found.

        If the request carries an idempotency key for a supported endpoint and
        a cached response exists, sends the cached response and returns True.
        Otherwise returns False so the caller proceeds with normal handling.

        Thread-safe. Uses module-level _idempotency_cache with TTL.

        Args:
            path: The request path (e.g., '/api/chat').

        Returns:
            True if a cached response was sent (caller should return early).
            False if no cached response -- caller should process normally.
        """
        if path not in _IDEMPOTENT_ENDPOINTS:
            return False
        idem_key = self.headers.get("X-Idempotency-Key") or ""
        if not idem_key:
            return False
        # Validate key format: 1-128 printable ASCII chars
        if len(idem_key) > 128 or not idem_key.isprintable():
            return False
        cached = _idempotency_get(idem_key)
        if cached is None:
            return False
        # Send cached response
        response_body: bytes = cached["response_body"]
        status_code: int = cached["status_code"]
        self.send_response(status_code)
        self.send_header("Content-Type", "application/json")
        self.send_header("X-Idempotent-Replay", "true")
        cors_origin = self._get_cors_origin()
        if cors_origin:
            self.send_header("Access-Control-Allow-Origin", cors_origin)
        self.send_header("Content-Length", str(len(response_body)))
        self.end_headers()
        self.wfile.write(response_body)
        return True

    def _store_idempotency(
        self, path: str, response_body: bytes, status_code: int
    ) -> None:
        """Store a response in the idempotency cache for the current request.

        Only stores if the request carried an X-Idempotency-Key header and the
        endpoint supports idempotency.

        Args:
            path: The request path.
            response_body: The serialized response bytes.
            status_code: The HTTP status code.
        """
        if path not in _IDEMPOTENT_ENDPOINTS:
            return
        idem_key = self.headers.get("X-Idempotency-Key") or ""
        if not idem_key or len(idem_key) > 128:
            return
        _idempotency_store(idem_key, response_body, status_code)

    def _check_admin_auth(self):
        """Check for admin API key via Authorization header or X-Admin-Key header.
        Uses hmac.compare_digest for timing-safe comparison to prevent
        timing side-channel attacks on the API key.
        SECURITY: Fails closed -- rejects if ADMIN_API_KEY is not configured.
        NOTE: Query parameter auth (?key=) is deliberately NOT supported because
        URL query strings are logged in server access logs, browser history,
        proxy logs, and Referer headers -- all of which leak the secret."""
        if not ADMIN_API_KEY:
            return False  # No key configured = reject (fail closed)
        import hmac

        # Method 1: Authorization: Bearer <key> header (API clients)
        auth = self.headers.get("Authorization") or ""
        if auth.startswith("Bearer "):
            key = auth[7:]
            if key and hmac.compare_digest(key, ADMIN_API_KEY):
                return True
        # Method 2: X-Admin-Key header (dashboard / browser access)
        admin_key_header = self.headers.get("X-Admin-Key") or ""
        if admin_key_header and hmac.compare_digest(admin_key_header, ADMIN_API_KEY):
            return True
        return False

    def _check_joveo_auth(self) -> bool:
        """Check if request is from an authenticated @joveo.com user or has valid API key.

        Three authentication paths (S58 hardened after security audit):
          1. X-Nova-Api-Key header (for embedded widgets on CG/GeoViz)
          2. Supabase JWT with @joveo.com email (HMAC-verified)
          3. Cookie-based session signed with SESSION_SIGNING_SECRET

        Returns:
            True if the request is authorized, False otherwise.
        """
        import base64
        import hmac as _hmac_mod
        import hashlib as _hashlib_mod

        # Path 1: API key bypass (for embedded widgets on CG/GeoViz)
        api_key = self.headers.get("X-Nova-Api-Key") or ""
        if api_key:
            allowed_keys = (os.environ.get("NOVA_API_KEYS") or "").split(",")
            if api_key.strip() in [k.strip() for k in allowed_keys if k.strip()]:
                return True

        # Path 2: Supabase JWT with @joveo.com email.  S58 FIX: previously the
        # payload was decoded without HMAC verification, so any attacker could
        # forge `Bearer x.<base64({"email":"x@joveo.com"})>.x` and bypass.  Now
        # we require a SUPABASE_JWT_SECRET env var and verify the signature.
        # If the secret is not configured, we FAIL CLOSED (reject the token)
        # rather than trusting it like the old code did.
        auth_header = self.headers.get("Authorization") or ""
        if auth_header.startswith("Bearer "):
            token = auth_header[7:]
            jwt_secret = (os.environ.get("SUPABASE_JWT_SECRET") or "").strip()
            if jwt_secret:
                try:
                    parts = token.split(".")
                    if len(parts) == 3:
                        header_b64, payload_b64, sig_b64 = parts
                        signing_input = f"{header_b64}.{payload_b64}".encode("ascii")
                        expected_sig = _hmac_mod.new(
                            jwt_secret.encode("utf-8"),
                            signing_input,
                            _hashlib_mod.sha256,
                        ).digest()
                        # Normalize base64url padding
                        actual_sig = base64.urlsafe_b64decode(
                            sig_b64 + "=" * (4 - len(sig_b64) % 4)
                        )
                        if _hmac_mod.compare_digest(expected_sig, actual_sig):
                            payload_padded = payload_b64 + "=" * (
                                4 - len(payload_b64) % 4
                            )
                            decoded = json.loads(
                                base64.urlsafe_b64decode(payload_padded)
                            )
                            # Verify token not expired (iat + exp claims)
                            _exp = decoded.get("exp")
                            if isinstance(_exp, (int, float)) and _exp < time.time():
                                logger.info("JWT expired (exp=%s)", _exp)
                            else:
                                email = (decoded.get("email") or "").lower().strip()
                                if email.endswith("@joveo.com"):
                                    return True
                except Exception as _jwt_err:
                    logger.debug("JWT verify failed: %s", _jwt_err)
            # If SUPABASE_JWT_SECRET is unset, refuse to validate unsigned
            # JWT claims (fail closed).

        # Path 3: Cookie-based session.  S58 FIX: previously any attacker
        # could set Cookie: nova_user_email=x@joveo.com and bypass.  Now,
        # when SESSION_SIGNING_SECRET is set, we require an HMAC-signed
        # cookie: email|hex(HMAC-SHA256(secret, email)).
        #
        # GRACEFUL MIGRATION: if SESSION_SIGNING_SECRET is unset OR
        # STRICT_AUTH is not enabled, we also accept the legacy unsigned
        # cookie and log a warning.  Operator flips STRICT_AUTH=1 after
        # verifying the signed-cookie flow works end-to-end.  This avoids
        # breaking every logged-in session the minute this deploys.
        cookie = self.headers.get("Cookie") or ""
        raw_cookie = _parse_cookie_value(cookie, "nova_user_email") or ""
        session_secret = (os.environ.get("SESSION_SIGNING_SECRET") or "").strip()
        strict_auth = (os.environ.get("STRICT_AUTH") or "").strip().lower() in (
            "1",
            "true",
            "yes",
            "on",
        )
        if raw_cookie:
            if session_secret and "|" in raw_cookie:
                try:
                    email_part, sig_part = raw_cookie.rsplit("|", 1)
                    expected = _hmac_mod.new(
                        session_secret.encode("utf-8"),
                        email_part.lower().strip().encode("utf-8"),
                        _hashlib_mod.sha256,
                    ).hexdigest()
                    if _hmac_mod.compare_digest(expected, sig_part.lower()):
                        if email_part.lower().strip().endswith("@joveo.com"):
                            return True
                except Exception as _cookie_err:
                    logger.debug("Cookie HMAC verify failed: %s", _cookie_err)
            if not strict_auth:
                email_lower = raw_cookie.lower().strip()
                if email_lower.endswith("@joveo.com"):
                    logger.warning(
                        "LEGACY-COOKIE-AUTH: unsigned nova_user_email "
                        "accepted for %s -- set STRICT_AUTH=1 after "
                        "cookie-issuance migration",
                        email_lower,
                    )
                    return True

        # Path 4: Allow embedded widgets from known Joveo product domains
        # This is the fallback when NOVA_API_KEYS is not configured yet.
        # Checks Origin/Referer header against known Joveo product URLs.
        # S48 FIX: ONLY external widget domains here -- NOT the main domain.
        # The main domain was allowing ANY visitor to bypass auth via Referer header.
        _joveo_widget_domains = {
            "cg-automation.onrender.com",
            "geoviz-3d.vercel.app",
            "geoviz.joveo.com",
        }
        origin = self.headers.get("Origin") or ""
        referer = self.headers.get("Referer") or ""
        for _jd in _joveo_widget_domains:
            if _jd in origin or _jd in referer:
                return True

        return False

    def _extract_authenticated_email(self) -> str:
        """Return the lowercase @joveo.com email if the request carries a verified
        session (HMAC-verified JWT or HMAC-signed cookie), else empty string.

        This mirrors the decoding logic in ``_check_joveo_auth`` but returns the
        identity instead of a bool so the rate limiter can key per-user.  Fails
        closed: if signatures do not verify, returns "".
        """
        import base64
        import hmac as _hmac_mod
        import hashlib as _hashlib_mod

        # Path A: HMAC-verified Supabase JWT
        auth_header = self.headers.get("Authorization") or ""
        if auth_header.startswith("Bearer "):
            token = auth_header[7:]
            jwt_secret = (os.environ.get("SUPABASE_JWT_SECRET") or "").strip()
            if jwt_secret:
                try:
                    parts = token.split(".")
                    if len(parts) == 3:
                        header_b64, payload_b64, sig_b64 = parts
                        signing_input = f"{header_b64}.{payload_b64}".encode("ascii")
                        expected_sig = _hmac_mod.new(
                            jwt_secret.encode("utf-8"),
                            signing_input,
                            _hashlib_mod.sha256,
                        ).digest()
                        actual_sig = base64.urlsafe_b64decode(
                            sig_b64 + "=" * (4 - len(sig_b64) % 4)
                        )
                        if _hmac_mod.compare_digest(expected_sig, actual_sig):
                            payload_padded = payload_b64 + "=" * (
                                4 - len(payload_b64) % 4
                            )
                            decoded = json.loads(
                                base64.urlsafe_b64decode(payload_padded)
                            )
                            _exp = decoded.get("exp")
                            if not (
                                isinstance(_exp, (int, float)) and _exp < time.time()
                            ):
                                email = (decoded.get("email") or "").lower().strip()
                                if email.endswith("@joveo.com"):
                                    return email
                except Exception as _jwt_err:
                    logger.debug("JWT email extract failed: %s", _jwt_err)

        # Path B: HMAC-signed (or legacy when STRICT_AUTH off) cookie
        cookie = self.headers.get("Cookie") or ""
        raw_cookie = _parse_cookie_value(cookie, "nova_user_email") or ""
        session_secret = (os.environ.get("SESSION_SIGNING_SECRET") or "").strip()
        strict_auth = (os.environ.get("STRICT_AUTH") or "").strip().lower() in (
            "1",
            "true",
            "yes",
            "on",
        )
        if raw_cookie:
            if session_secret and "|" in raw_cookie:
                try:
                    email_part, sig_part = raw_cookie.rsplit("|", 1)
                    expected = _hmac_mod.new(
                        session_secret.encode("utf-8"),
                        email_part.lower().strip().encode("utf-8"),
                        _hashlib_mod.sha256,
                    ).hexdigest()
                    if _hmac_mod.compare_digest(expected, sig_part.lower()):
                        email_lower = email_part.lower().strip()
                        if email_lower.endswith("@joveo.com"):
                            return email_lower
                except Exception as _cookie_err:
                    logger.debug("Cookie email extract failed: %s", _cookie_err)
            if not strict_auth:
                email_lower = raw_cookie.lower().strip()
                if email_lower.endswith("@joveo.com"):
                    return email_lower

        return ""

    def _check_rate_limit(self):
        """Tiered rate limiting: API key tier limits take precedence over per-IP limits.

        If the request carries a valid ``Authorization: Bearer <key>`` whose key
        exists in ``_api_keys_store`` and is not revoked, the per-minute and
        per-day limits from the key's tier are used instead of the default
        per-IP limits.  Otherwise falls back to original IP-based limiting.

        Thread-safe. Cleans up stale entries to prevent unbounded memory growth.
        """
        now = time.time()

        # ── Check for API-key-based tier limits ──
        auth_header = self.headers.get("Authorization") or ""
        api_key = None
        if auth_header.startswith("Bearer "):
            api_key = auth_header[7:]

        if api_key:
            with _api_keys_lock:
                key_entry = _api_keys_store.get(api_key)
                if key_entry and not key_entry.get("revoked"):
                    tier_name = key_entry.get("tier", "free")
                    tier_limits = API_KEY_TIERS.get(tier_name, API_KEY_TIERS["free"])

                    # Per-minute check
                    key_entry["usage_minute"] = [
                        t for t in key_entry.get("usage_minute") or [] if now - t < 60
                    ]
                    if len(key_entry["usage_minute"]) >= tier_limits["rpm"]:
                        return False

                    # Per-day check
                    key_entry["usage_day"] = [
                        t for t in key_entry.get("usage_day") or [] if now - t < 86400
                    ]
                    if len(key_entry["usage_day"]) >= tier_limits["rpd"]:
                        return False

                    key_entry["usage_minute"].append(now)
                    key_entry["usage_day"].append(now)
                    return True
                # If key is invalid/revoked, fall through to IP-based limiting

        # ── Fallback: per-IP rate limiting (original behavior) ──
        # Also layer per-user limiting on top for authenticated @joveo.com
        # users, because the entire Joveo office shares one NAT egress IP
        # and one heavy-testing engineer must not starve colleagues out of
        # the shared per-IP bucket.
        client_ip = self._get_client_ip()
        user_email = self._extract_authenticated_email()  # "" if not authed

        with _rate_limit_lock:
            # Purge expired timestamps for this IP
            _rate_limit_store[client_ip] = [
                t for t in _rate_limit_store[client_ip] if now - t < _RATE_LIMIT_WINDOW
            ]
            if len(_rate_limit_store[client_ip]) >= _RATE_LIMIT_MAX:
                logger.info(
                    "RATE-LIMIT (per-IP) exceeded: ip=%s count=%d cap=%d",
                    client_ip,
                    len(_rate_limit_store[client_ip]),
                    _RATE_LIMIT_MAX,
                )
                return False

            # Per-user check (only for authenticated @joveo.com users)
            if user_email:
                _user_rate_limit_store[user_email] = [
                    t
                    for t in _user_rate_limit_store[user_email]
                    if now - t < _RATE_LIMIT_WINDOW
                ]
                if len(_user_rate_limit_store[user_email]) >= _USER_RATE_LIMIT_MAX:
                    logger.info(
                        "RATE-LIMIT (per-user) exceeded: email=%s count=%d cap=%d "
                        "(ip=%s ok at %d/%d)",
                        user_email,
                        len(_user_rate_limit_store[user_email]),
                        _USER_RATE_LIMIT_MAX,
                        client_ip,
                        len(_rate_limit_store[client_ip]),
                        _RATE_LIMIT_MAX,
                    )
                    return False

            # Both buckets (IP, and user if authed) have capacity -- record.
            _rate_limit_store[client_ip].append(now)
            if user_email:
                _user_rate_limit_store[user_email].append(now)

            # Periodic cleanup: evict stale IPs and users with no recent requests.
            # Run cleanup when either store exceeds 500 entries to stay lean.
            if len(_rate_limit_store) > 500:
                stale = [
                    ip
                    for ip, ts in list(_rate_limit_store.items())
                    if not ts or now - max(ts) > _RATE_LIMIT_WINDOW * 2
                ]
                for ip in stale:
                    _rate_limit_store.pop(ip, None)
            if len(_user_rate_limit_store) > 500:
                stale_users = [
                    email
                    for email, ts in list(_user_rate_limit_store.items())
                    if not ts or now - max(ts) > _RATE_LIMIT_WINDOW * 2
                ]
                for email in stale_users:
                    _user_rate_limit_store.pop(email, None)
        return True

    def _check_global_chat_rate_limit(self):
        """Global rate limit across all IPs for /api/chat to prevent distributed cost abuse.

        Thread-safe. Returns True if request is allowed, False if limit exceeded.
        """
        now = time.time()
        with _global_chat_lock:
            # Purge expired timestamps
            _global_chat_timestamps[:] = [
                t for t in _global_chat_timestamps if now - t < _RATE_LIMIT_WINDOW
            ]
            if len(_global_chat_timestamps) >= _GLOBAL_CHAT_RATE_LIMIT_MAX:
                return False
            _global_chat_timestamps.append(now)
        return True

    # ══════════════════════════════════════════════════════════════════════════
    # WEBSOCKET CHAT HANDLER
    # ══════════════════════════════════════════════════════════════════════════
    def _handle_ws_chat(self) -> None:
        """Handle WebSocket upgrade for /ws/chat.

        Performs the WebSocket handshake, then enters a read loop.
        When the client sends a chat message (JSON), streams the Nova
        response back as individual WebSocket text frames -- same JSON
        event format as the SSE endpoint so the client parser is shared.

        S58 FIX: Previously this endpoint had ZERO authentication --
        anyone on the internet could connect and exhaust Anthropic
        credits. Now it enforces the same @joveo.com auth that SSE uses
        AND validates Origin against an allowlist to block CSWSH.
        """
        # ── Auth check BEFORE handshake so attackers cannot even upgrade ──
        if not self._check_joveo_auth():
            logger.warning(
                "WS /ws/chat: auth rejected (client=%s, origin=%s)",
                self._get_client_ip(),
                self.headers.get("Origin") or "",
            )
            self.send_error(401, "Authentication required")
            return

        # ── Origin allowlist: block CSWSH from arbitrary third-party sites ──
        origin = (self.headers.get("Origin") or "").lower()
        _allowed_ws_origins = {
            "https://media-plan-generator.onrender.com",
            "https://nova.joveo.com",
            "https://cg-automation.onrender.com",
            "https://geoviz-3d.vercel.app",
            "https://geoviz.joveo.com",
            "http://localhost:5000",
            "http://127.0.0.1:5000",
        }
        if origin and origin not in _allowed_ws_origins:
            logger.warning("WS /ws/chat: Origin rejected: %s", origin)
            self.send_error(403, "Origin not allowed")
            return

        # ── Per-connection rate limit: cap concurrent WS per IP ──
        if not self._check_rate_limit():
            logger.warning("WS /ws/chat: rate-limited for %s", self._get_client_ip())
            self.send_error(429, "Too many requests")
            return

        try:
            from websocket_handler import ws_handshake
        except ImportError:
            logger.error("websocket_handler module not available", exc_info=True)
            self.send_error(500, "WebSocket support unavailable")
            return

        ws_conn = ws_handshake(self)
        if ws_conn is None:
            return  # Handshake failed, error already sent

        logger.info("WebSocket chat connection opened: %s", ws_conn.connection_id)

        try:
            while not ws_conn.closed:
                raw_msg = ws_conn.recv(timeout=120.0)
                if raw_msg is None:
                    break  # Connection closed or timed out

                # Parse the incoming message
                try:
                    data = json.loads(raw_msg)
                except (json.JSONDecodeError, ValueError):
                    ws_conn.send_json({"error": "Invalid JSON", "done": True})
                    continue

                msg_type = data.get("type", "chat")

                if msg_type == "ping":
                    ws_conn.send_json({"type": "pong"})
                    continue

                if msg_type == "stop":
                    # Cancel active stream for this conversation
                    conv_id = data.get("conversation_id") or ""
                    if conv_id:
                        _cancel_stream(conv_id)
                    ws_conn.send_json({"type": "stop_ack", "conversation_id": conv_id})
                    continue

                if msg_type != "chat":
                    ws_conn.send_json(
                        {"error": f"Unknown message type: {msg_type}", "done": True}
                    )
                    continue

                # ── Chat message processing ──
                # Validate
                chat_msg = data.get("message") or ""
                if not chat_msg:
                    ws_conn.send_json({"error": "Empty message", "done": True})
                    continue

                if isinstance(chat_msg, str) and len(chat_msg) > 10_000:
                    ws_conn.send_json(
                        {
                            "error": "Message exceeds 10,000 character limit",
                            "done": True,
                        }
                    )
                    continue

                chat_history = data.get("history") or []
                if isinstance(chat_history, list) and len(chat_history) > 50:
                    ws_conn.send_json(
                        {
                            "error": "Chat history exceeds 50 messages limit",
                            "done": True,
                        }
                    )
                    continue

                # Sanitize and deduplicate repetitive input
                data["message"] = _sanitize_chat_input(data.get("message"))
                data["message"] = _dedup_repetitive_input(data["message"])

                conversation_id = data.get("conversation_id") or str(uuid.uuid4())
                ws_conn.conversation_id = conversation_id

                # Session token verification
                _conv_session_token = ""
                client_session_token = data.get("session_token") or ""
                try:
                    from nova_persistence import (
                        get_or_create_conversation,
                        verify_conversation_token,
                    )

                    conv_row = get_or_create_conversation(conversation_id)
                    if conv_row:
                        _conv_session_token = conv_row.get("session_token") or ""
                        if (
                            _conv_session_token
                            and client_session_token
                            and not verify_conversation_token(
                                conversation_id, client_session_token
                            )
                        ):
                            ws_conn.send_json(
                                {
                                    "error": "Unauthorized: invalid session token",
                                    "done": True,
                                }
                            )
                            continue
                except ImportError:
                    pass
                except Exception as _auth_err:
                    logger.error(
                        "WS conversation auth check failed: %s",
                        _auth_err,
                        exc_info=True,
                    )

                cancel_event = _register_stream(conversation_id)

                try:
                    # Enrichment
                    _enrich_chat_context(data, data.get("message") or "")

                    # Check cancellation
                    if cancel_event.is_set():
                        ws_conn.send_json(
                            {
                                "token": "",
                                "done": True,
                                "cancelled": True,
                                "full_response": "Your request was cancelled.",
                            }
                        )
                        continue

                    # Stream Nova response
                    from nova import handle_chat_request_stream

                    full_response = ""
                    sources: list = []
                    confidence = None  # S47: None means "don't show badge"
                    tools_used: list = []
                    model_used = ""
                    message_id = str(uuid.uuid4())
                    _token_usage: dict = {}
                    # S62: Transparency panel fields forwarded from done chunk.
                    _ws_fast_path = ""
                    _ws_kb_files_queried: list = []
                    _ws_timing_ms = None

                    # Send initial status
                    ws_conn.send_json({"type": "status", "status": "Thinking..."})

                    # Keepalive thread for WebSocket
                    _ws_keepalive_stop = threading.Event()
                    _ws_last_event = time.time()

                    def _ws_keepalive() -> None:
                        nonlocal _ws_last_event
                        while not _ws_keepalive_stop.is_set():
                            _ws_keepalive_stop.wait(15)
                            if _ws_keepalive_stop.is_set():
                                break
                            if time.time() - _ws_last_event >= 15:
                                if not ws_conn.send_json({"keepalive": True}):
                                    break
                                _ws_last_event = time.time()

                    _ka_t = threading.Thread(
                        target=_ws_keepalive, daemon=True, name="ws-keepalive"
                    )
                    _ka_t.start()

                    try:
                        for chunk in handle_chat_request_stream(
                            data, cancel_event=cancel_event
                        ):
                            if cancel_event.is_set():
                                ws_conn.send_json(
                                    {
                                        "token": "",
                                        "done": True,
                                        "cancelled": True,
                                        "full_response": full_response,
                                        "sources": sources,
                                    }
                                )
                                break

                            if chunk.get("done"):
                                full_response = (
                                    chunk.get("full_response") or full_response
                                )
                                sources = chunk.get("sources") or sources
                                # S47: Preserve None confidence (means "don't show badge")
                                _ws_chunk_conf = chunk.get("confidence")
                                if _ws_chunk_conf is not None:
                                    confidence = _ws_chunk_conf
                                tools_used = chunk.get("tools_used") or tools_used
                                model_used = chunk.get("llm_provider") or model_used
                                _token_usage = chunk.get("token_usage") or {}
                                # S62: Pull transparency panel fields from
                                # handle_chat_request_stream's done chunk so
                                # the final WS event can carry them to the
                                # frontend (fast_path + kb_files_queried +
                                # timing_ms were being dropped here).
                                _ws_fast_path = chunk.get("fast_path") or ""
                                _ws_kb_files_queried = (
                                    chunk.get("kb_files_queried") or []
                                )
                                _ws_timing_ms = chunk.get("timing_ms")
                            elif chunk.get("type") in ("tool_start", "tool_complete"):
                                # S18: Real-time tool status via WebSocket
                                if not ws_conn.send_json(
                                    {
                                        "type": chunk["type"],
                                        "tool": chunk.get("tool", ""),
                                        "label": chunk.get("label", ""),
                                        "done": False,
                                    }
                                ):
                                    break
                                _ws_last_event = time.time()
                            elif chunk.get("status"):
                                if not ws_conn.send_json(
                                    {"status": chunk["status"], "done": False}
                                ):
                                    break
                                _ws_last_event = time.time()
                            else:
                                token = chunk.get("token") or ""
                                if token:
                                    full_response += token
                                    if not ws_conn.send_json(
                                        {"token": token, "done": False}
                                    ):
                                        break
                                    _ws_last_event = time.time()
                    finally:
                        _ws_keepalive_stop.set()

                    # Follow-up suggestions
                    _ws_followups: list[str] = []
                    try:
                        _ws_user_msg = data.get("message") or ""
                        if full_response and _ws_user_msg:
                            _ws_followups = _generate_followup_questions(
                                full_response, _ws_user_msg
                            )
                    except Exception as _wsfq:
                        logger.error(
                            "WS follow-up gen failed: %s", _wsfq, exc_info=True
                        )

                    # Final event with metadata
                    ws_conn.send_json(
                        {
                            "token": "",
                            "done": True,
                            "full_response": full_response,
                            "sources": sources,
                            "confidence": confidence,
                            "tools_used": tools_used,
                            "message_id": message_id,
                            "conversation_id": conversation_id,
                            "suggested_followups": _ws_followups,
                            "session_token": _conv_session_token,
                            "token_usage": _token_usage,
                            # S62: Transparency panel payload for "Why this
                            # answer?" footer. Forwarded from the Nova stream
                            # done chunk so the WS client sees real values
                            # instead of None / [] on fast-path responses.
                            "fast_path": _ws_fast_path,
                            "kb_files_queried": _ws_kb_files_queried,
                            "timing_ms": _ws_timing_ms,
                        }
                    )

                    # Persist to Supabase asynchronously
                    _ws_user_msg = data.get("message") or ""
                    if conversation_id and _ws_user_msg:
                        threading.Thread(
                            target=_save_conversation_turn,
                            args=(
                                conversation_id,
                                _ws_user_msg,
                                full_response,
                                model_used,
                                sources,
                                confidence,
                            ),
                            daemon=True,
                            name="supabase-ws-persist",
                        ).start()

                except Exception as ws_stream_err:
                    logger.error(
                        "WebSocket streaming error: %s", ws_stream_err, exc_info=True
                    )
                    ws_conn.send_json(
                        {
                            "token": "",
                            "done": True,
                            "error": "An error occurred processing your request.",
                        }
                    )
                finally:
                    _unregister_stream(conversation_id)

        except Exception as ws_err:
            logger.error("WebSocket handler error: %s", ws_err, exc_info=True)
        finally:
            ws_conn.close()
            logger.info("WebSocket chat connection closed: %s", ws_conn.connection_id)

    def do_GET(self):
        # ── Request ID generation (Feature 6) ──
        try:
            from monitoring import generate_request_id as _gen_rid

            self._request_id = _gen_rid()
        except Exception:
            self._request_id = uuid.uuid4().hex[:12]
        # ── Set request context for structured logging ──
        try:
            from monitoring import set_request_id as _set_rid

            _set_rid(self._request_id)
        except ImportError:
            pass
        _req_start = time.time()
        if _metrics:
            _metrics.enter_request()

        # ── Sentry Performance: wrap entire GET in a transaction (SDK 2.x) ──
        parsed_for_sentry = urlparse(self.path)
        _sentry_txn_cm, self._sentry_span = _sentry_start_txn(
            "http.server", f"GET {parsed_for_sentry.path}"
        )

        self._response_status: int = 200
        try:
            with _sentry_txn_cm:
                self._handle_GET()
        except Exception as _exc:
            # Sentry captures unhandled exceptions automatically; log + re-raise
            logger.exception("Unhandled exception in do_GET: %s", _exc)
            self._response_status = 500
            # Slack alert for unhandled GET errors
            if _slack_alerts_available:
                try:
                    _slack_error_alert(
                        _exc,
                        context="Unhandled exception in do_GET",
                        endpoint=urlparse(self.path).path,
                    )
                except Exception:
                    pass
            try:
                self._send_error(
                    "Internal server error",
                    "INTERNAL_ERROR",
                    500,
                )
            except Exception:
                pass
        finally:
            _latency = (time.time() - _req_start) * 1000
            if _metrics:
                _metrics.exit_request()
                parsed = urlparse(self.path)
                _metrics.record_request(
                    parsed.path, "GET", self._response_status, _latency
                )
            # ── Slow endpoint profiling ──
            _record_slow_endpoint(urlparse(self.path).path, "GET", _latency)
            # ── SLO tracker + anomaly detector (S17) ──
            try:
                from slo_config import get_slo_tracker

                get_slo_tracker().record_from_endpoint(
                    urlparse(self.path).path, _latency, self._response_status
                )
            except Exception:
                pass
            try:
                from anomaly_detector import record_metric, METRIC_REQUEST_LATENCY

                record_metric(METRIC_REQUEST_LATENCY, _latency)
            except Exception:
                pass
            # ── PostHog: Track page views for HTML template routes ──
            if _posthog_available and self._response_status == 200:
                _ph_path = urlparse(self.path).path.rstrip("/") or "/"
                # Only track page views, not API/static/health endpoints
                _PH_SKIP_PREFIXES = (
                    "/api/",
                    "/static/",
                    "/health",
                    "/ready",
                    "/fragment/",
                )
                _PH_SKIP_EXACT = ("/robots.txt", "/sitemap.xml", "/favicon.ico")
                if (
                    not any(_ph_path.startswith(p) for p in _PH_SKIP_PREFIXES)
                    and _ph_path not in _PH_SKIP_EXACT
                ):
                    try:
                        self._ph_track_page(_ph_path, _ph_path.strip("/") or "hub")
                    except Exception:
                        pass  # Never let analytics break the request
            # ── Clear request context for thread reuse ──
            try:
                from monitoring import clear_request_context as _clear_ctx

                _clear_ctx()
            except ImportError:
                pass

    def _handle_GET(self):
        parsed = urlparse(self.path)
        path = parsed.path
        # ── API Versioning: strip /v1 prefix (Feature 4) ──
        if path.startswith("/v1/"):
            path = path[3:]

        # ── WebSocket Upgrade: intercept before normal GET routing ──
        upgrade_header = (self.headers.get("Upgrade") or "").lower()
        if upgrade_header == "websocket" and path == "/ws/chat":
            self._handle_ws_chat()
            return

        # ── API Key Authentication (Phase 6) ──
        if _auth_module_loaded:
            try:
                from auth import authenticate

                auth_header = self.headers.get("Authorization") or ""
                query_params = urllib.parse.parse_qs(
                    urllib.parse.urlparse(self.path).query
                )
                query_key = query_params.get("api_key", [None])[0]
                auth_result = authenticate(path, auth_header, query_key)
                if not auth_result["authenticated"]:
                    self._send_json(
                        {
                            "error": "Authentication required",
                            "reason": auth_result["reason"],
                        },
                        status_code=401,
                    )
                    return
            except Exception:
                logger.error("Auth check failed unexpectedly", exc_info=True)
                self._send_json({"error": "Internal auth error"}, status_code=500)
                return
        elif path.startswith("/api/admin"):
            # Auth module not loaded -- block admin endpoints
            self._send_json(
                {"error": "Auth module unavailable -- admin access denied"},
                status_code=503,
            )
            return

        # ── Audit logging for API access ──
        try:
            from audit_logger import log_event

            if path.startswith("/api/"):
                log_event(
                    "api.access", resource=path, ip_address=self.client_address[0]
                )
        except ImportError:
            pass

        # ── Extracted route modules (health, admin, chat GET, pages, campaign, pricing) ──
        from routes.health import handle_health_routes
        from routes.admin import handle_admin_routes
        from routes.chat import handle_chat_get_routes
        from routes.pages import handle_page_routes
        from routes.campaign import handle_campaign_get_routes
        from routes.pricing import handle_pricing_get_routes
        from routes.canvas import handle_canvas_get_routes
        from routes.export import handle_export_get_routes
        from routes.slotops import handle_slotops_get_routes

        if handle_health_routes(self, path, parsed):
            return
        if handle_admin_routes(self, path, parsed):
            return
        if handle_chat_get_routes(self, path, parsed):
            return
        if handle_page_routes(self, path, parsed):
            return
        if handle_campaign_get_routes(self, path, parsed):
            return
        if handle_pricing_get_routes(self, path, parsed):
            return
        if handle_canvas_get_routes(self, path, parsed):
            return
        if handle_export_get_routes(self, path, parsed):
            return
        if handle_slotops_get_routes(self, path, parsed):
            return
        # NOTE: /plan/shared/<id>, page routes (platform, health-dashboard, etc.)
        # are now handled by routes/pages.py and routes/campaign.py above.
        # ── Fragment Endpoint (serves HTML fragments for platform SPA) ──
        if path.startswith("/fragment/"):
            fragment_name = path[len("/fragment/") :].strip("/")
            self._serve_fragment(fragment_name)

        # ── Public Config Endpoint (auth gate + NovaAuth initialization) ──
        elif path == "/api/config":
            _sb_url = os.environ.get("SUPABASE_URL") or ""
            _sb_anon = os.environ.get("SUPABASE_ANON_KEY") or ""
            _auth_enabled = bool(_sb_url and _sb_anon)
            if not _auth_enabled:
                logger.warning(
                    "[Auth Config] SUPABASE_URL or SUPABASE_ANON_KEY not set -- "
                    "Google OAuth will not work. Set both env vars on Render."
                )
            self._send_json(
                {
                    "auth_enabled": _auth_enabled,
                    "supabase_url": _sb_url if _auth_enabled else "",
                    "supabase_anon_key": _sb_anon if _auth_enabled else "",
                    "allowed_domains": ["joveo.com"],
                    "provider": "google",
                }
            )

        elif path == "/api/csrf-token":
            # Cookie-based double-submit CSRF: generate token, set as cookie
            # AND return in JSON body.  JS stores the body copy and sends it
            # back via X-CSRF-Token header; browser sends the cookie copy
            # automatically.  Server compares the two on POST.
            token = _generate_csrf_token()
            is_https = (self.headers.get("X-Forwarded-Proto") or "").lower() == "https"
            cookie_val = _build_csrf_cookie(token, secure=is_https)
            body = json.dumps({"token": token}).encode("utf-8")
            self.send_response(200)
            self.send_header("Content-Type", "application/json")
            self.send_header("Set-Cookie", cookie_val)
            # C1 fix: Set a separate HttpOnly session cookie for IDOR prevention
            session_cookie_val = _build_session_cookie(token, secure=is_https)
            self.send_header("Set-Cookie", session_cookie_val)
            cors_origin = self._get_cors_origin()
            if cors_origin:
                self.send_header("Access-Control-Allow-Origin", cors_origin)
                self.send_header("Access-Control-Allow-Credentials", "true")
            self.send_header("Content-Length", str(len(body)))
            self.end_headers()
            self.wfile.write(body)

        elif path == "/api/auth/me":
            # ── Auth: Return current user info (from Supabase JWT) ──
            # Non-breaking: returns null user if not authenticated
            auth_header = self.headers.get("Authorization") or ""
            if auth_header.startswith("Bearer "):
                _auth_token = auth_header[7:]
                # Validate via Supabase REST API
                _sb_url = os.environ.get("SUPABASE_URL") or ""
                _sb_key = os.environ.get("SUPABASE_ANON_KEY") or ""
                if _sb_url and _sb_key:
                    try:
                        _auth_req = urllib.request.Request(
                            f"{_sb_url.rstrip('/')}/auth/v1/user",
                            headers={
                                "apikey": _sb_key,
                                "Authorization": f"Bearer {_auth_token}",
                            },
                        )
                        with urllib.request.urlopen(_auth_req, timeout=5) as _resp:
                            _user_data = json.loads(_resp.read().decode())
                            _user_email = _user_data.get("email") or ""
                            _user_name = (_user_data.get("user_metadata") or {}).get(
                                "full_name", ""
                            )
                            _user_avatar = (_user_data.get("user_metadata") or {}).get(
                                "avatar_url", ""
                            )

                            # S46: Server-side domain enforcement
                            if not _user_email.lower().strip().endswith("@joveo.com"):
                                self._send_json(
                                    {
                                        "authenticated": False,
                                        "error": "Only @joveo.com accounts are allowed",
                                    },
                                    status_code=403,
                                )
                                return

                            # S40: Log successful auth to nova_login_log (best-effort)
                            try:
                                import hashlib as _hl

                                _client_ip = (
                                    self.headers.get("X-Forwarded-For", "")
                                    .split(",")[0]
                                    .strip()
                                    or self.client_address[0]
                                )
                                _ip_hash = _hl.sha256(_client_ip.encode()).hexdigest()[
                                    :16
                                ]
                                _ua = (self.headers.get("User-Agent") or "")[:200]
                                _supabase_rest(
                                    "nova_login_log",
                                    method="POST",
                                    payload={
                                        "user_email": _user_email,
                                        "user_name": _user_name,
                                        "provider": "google",
                                        "ip_hash": _ip_hash,
                                        "user_agent": _ua,
                                        "product": "nova",
                                    },
                                )
                            except Exception as _login_log_err:
                                logger.debug(
                                    "nova_login_log write failed: %s", _login_log_err
                                )

                            self._send_json(
                                {
                                    "authenticated": True,
                                    "user": {
                                        "id": _user_data.get("id", ""),
                                        "email": _user_email,
                                        "name": _user_name,
                                        "avatar_url": _user_avatar,
                                    },
                                }
                            )
                            return
                    except Exception:
                        pass
            # Not authenticated -- return gracefully (non-breaking)
            self._send_json({"authenticated": False, "user": None})

        elif path == "/api/slack/status":
            # ── Slack Bot Status Endpoint (admin-protected) ──
            # H4 fix: Only expose healthy/unhealthy status -- no internal
            # configuration, key presence, URLs, or exception details.
            if not self._check_admin_auth():
                self._send_json({"error": "Unauthorized"}, status_code=401)
                return
            _slack_bot_token = os.environ.get("SLACK_BOT_TOKEN") or ""
            _slack_signing_secret = os.environ.get("SLACK_SIGNING_SECRET") or ""
            _is_healthy = bool(_slack_bot_token and _slack_signing_secret)
            if _is_healthy:
                try:
                    from nova_slack import get_nova_bot

                    bot = get_nova_bot()
                    _is_healthy = bool(bot.bot_user_id)
                except Exception as e:
                    logger.error("Slack bot health check failed: %s", e, exc_info=True)
                    _is_healthy = False
            if _is_healthy:
                self._send_json({"status": "healthy"})
            else:
                self._send_json(
                    {"status": "unhealthy", "error": "service unavailable"},
                    status_code=503,
                )
        elif path.startswith("/google") and path.endswith(".html") and len(path) < 60:
            # Google Search Console verification file
            _static = os.path.join(os.path.dirname(os.path.abspath(__file__)), "static")
            verify_file = os.path.join(_static, path.lstrip("/"))
            if os.path.isfile(verify_file):
                with open(verify_file, "r") as f:
                    content = f.read()
                self.send_response(200)
                self.send_header("Content-Type", "text/html")
                self.end_headers()
                self.wfile.write(content.encode())
            else:
                self.send_error(404)
            return

        elif path == "/favicon.ico":
            # Redirect to SVG favicon
            self.send_response(301)
            self.send_header("Location", "/static/favicon.svg")
            self.send_header("Cache-Control", "public, max-age=604800")
            self.end_headers()
        elif path == "/robots.txt":
            robots_content = (
                "User-agent: *\n"
                "Allow: /\n"
                "Disallow: /api/\n"
                "Disallow: /admin/\n"
                "Disallow: /health\n"
                "Disallow: /ready\n"
                "Disallow: /dashboard\n"
                "\n"
                "# AI crawlers welcome (GEO/AEO optimization)\n"
                "User-agent: GPTBot\n"
                "Allow: /\n"
                "Disallow: /api/\n"
                "\n"
                "User-agent: ChatGPT-User\n"
                "Allow: /\n"
                "Disallow: /api/\n"
                "\n"
                "User-agent: Claude-Web\n"
                "Allow: /\n"
                "Disallow: /api/\n"
                "\n"
                "User-agent: Google-Extended\n"
                "Allow: /\n"
                "Disallow: /api/\n"
                "\n"
                "Sitemap: https://media-plan-generator.onrender.com/sitemap.xml\n"
            )
            self.send_response(200)
            self.send_header("Content-Type", "text/plain; charset=utf-8")
            self.send_header("Cache-Control", "public, max-age=86400")
            self.end_headers()
            self.wfile.write(robots_content.encode("utf-8"))
        elif path == "/sitemap.xml":
            _today = datetime.date.today().isoformat()
            _BASE = "https://media-plan-generator.onrender.com"
            # All public product pages with SEO priority
            _SITEMAP_PAGES: list[tuple[str, str, str]] = [
                ("/", "1.0", "weekly"),
                ("/platform", "0.9", "weekly"),
                ("/platform/plan", "0.85", "weekly"),
                ("/platform/intelligence", "0.85", "weekly"),
                ("/platform/compliance", "0.85", "weekly"),
                ("/platform/nova", "0.8", "weekly"),
                ("/nova", "0.9", "weekly"),
                ("/media-plan", "0.9", "weekly"),
                ("/competitive-intel", "0.8", "weekly"),
                ("/market-pulse", "0.8", "weekly"),
                ("/market-intel", "0.8", "weekly"),
                ("/vendor-iq", "0.8", "weekly"),
                ("/talent-heatmap", "0.8", "weekly"),
                ("/hire-signal", "0.8", "weekly"),
                ("/payscale-sync", "0.8", "weekly"),
                ("/skill-target", "0.8", "weekly"),
                ("/applyflow", "0.7", "weekly"),
                ("/compliance-guard", "0.8", "weekly"),
                ("/audit", "0.7", "weekly"),
                ("/budget-simulator", "0.7", "weekly"),
                ("/roi-calculator", "0.7", "weekly"),
                ("/performance-tracker", "0.7", "weekly"),
                ("/post-campaign", "0.7", "weekly"),
                ("/quick-plan", "0.7", "weekly"),
                ("/quick-brief", "0.7", "weekly"),
                ("/social-plan", "0.7", "weekly"),
                ("/creative-ai", "0.7", "weekly"),
                ("/ab-testing", "0.7", "weekly"),
                ("/pricing", "0.6", "monthly"),
                ("/privacy", "0.3", "monthly"),
                ("/terms", "0.3", "monthly"),
            ]
            _url_entries = ""
            for _pg_path, _pg_prio, _pg_freq in _SITEMAP_PAGES:
                _url_entries += (
                    f"  <url>\n"
                    f"    <loc>{_BASE}{_pg_path}</loc>\n"
                    f"    <lastmod>{_today}</lastmod>\n"
                    f"    <changefreq>{_pg_freq}</changefreq>\n"
                    f"    <priority>{_pg_prio}</priority>\n"
                    f"  </url>\n"
                )
            sitemap_content = (
                '<?xml version="1.0" encoding="UTF-8"?>\n'
                '<urlset xmlns="http://www.sitemaps.org/schemas/sitemap/0.9">\n'
                f"{_url_entries}"
                "</urlset>\n"
            )
            self.send_response(200)
            self.send_header("Content-Type", "application/xml; charset=utf-8")
            self.send_header("Cache-Control", "public, max-age=86400")
            self.end_headers()
            self.wfile.write(sitemap_content.encode("utf-8"))
        # ── Plan Templates Marketplace ──
        elif path == "/api/templates":
            try:
                from plan_templates import list_templates as _list_templates

                query_params = urllib.parse.parse_qs(parsed.query)
                category = query_params.get("category", [None])[0]
                role_family = query_params.get("role_family", [None])[0]
                templates = _list_templates(category=category, role_family=role_family)
                self._send_json({"total": len(templates), "templates": templates})
            except ImportError:
                self._send_json(
                    {"error": "Plan templates module not available"}, status_code=503
                )
            except (ValueError, TypeError, KeyError) as e:
                logger.error(f"Templates list error: {e}", exc_info=True)
                self._send_json(
                    {"error": f"Failed to list templates: {e}"}, status_code=500
                )
        elif path.startswith("/api/templates/") and not path.endswith("/fork"):
            try:
                from plan_templates import get_template as _get_template

                template_id = path[len("/api/templates/") :]
                tmpl = _get_template(template_id)
                if tmpl is None:
                    self._send_json(
                        {"error": f"Template '{template_id}' not found"},
                        status_code=404,
                    )
                else:
                    self._send_json(tmpl)
            except ImportError:
                self._send_json(
                    {"error": "Plan templates module not available"}, status_code=503
                )
            except (ValueError, TypeError, KeyError) as e:
                logger.error(f"Template get error: {e}", exc_info=True)
                self._send_json(
                    {"error": f"Failed to get template: {e}"}, status_code=500
                )
        elif path == "/api/channels":
            db = load_channels_db()
            # Inject the full industry options list for frontend consumption
            db["industry_options"] = [
                {"value": "healthcare_medical", "label": "Healthcare & Medical"},
                {
                    "value": "blue_collar_trades",
                    "label": "Blue Collar / Skilled Trades",
                },
                {"value": "maritime_marine", "label": "Maritime & Marine"},
                {"value": "military_recruitment", "label": "Military Recruitment"},
                {"value": "tech_engineering", "label": "Technology & Engineering"},
                {"value": "general_entry_level", "label": "General / Entry-Level"},
                {"value": "legal_services", "label": "Legal Services"},
                {"value": "finance_banking", "label": "Finance & Banking"},
                {"value": "mental_health", "label": "Mental Health & Behavioral"},
                {"value": "retail_consumer", "label": "Retail & Consumer"},
                {"value": "aerospace_defense", "label": "Aerospace & Defense"},
                {"value": "pharma_biotech", "label": "Pharma & Biotech"},
                {"value": "energy_utilities", "label": "Energy & Utilities"},
                {"value": "insurance", "label": "Insurance"},
                {"value": "telecommunications", "label": "Telecommunications"},
                {"value": "automotive", "label": "Automotive & Manufacturing"},
                {"value": "food_beverage", "label": "Food & Beverage"},
                {
                    "value": "logistics_supply_chain",
                    "label": "Logistics & Supply Chain",
                },
                {"value": "hospitality_travel", "label": "Hospitality & Travel"},
                {"value": "media_entertainment", "label": "Media & Entertainment"},
                {
                    "value": "construction_real_estate",
                    "label": "Construction & Real Estate",
                },
                {"value": "education", "label": "Education"},
            ]
            self._send_json(db)
        elif path == "/api/requests":
            if not self._check_admin_auth():
                self._send_json({"error": "Unauthorized"}, status_code=401)
                return
            log = load_request_log()
            # Add download URLs for entries that have doc_filename
            enriched_log = []
            for entry in log:
                e = dict(entry)
                if e.get("doc_filename"):
                    e["doc_download_url"] = f"/api/documents/{e['doc_filename']}"
                enriched_log.append(e)
            self.send_response(200)
            self.send_header("Content-Type", "application/json")
            self.end_headers()
            response = {
                "total_requests": len(enriched_log),
                "requests": enriched_log[-100:],  # Last 100 entries
            }
            self.wfile.write(json.dumps(response, indent=2, default=str).encode())
        # NOTE: /dashboard and /observability now handled by routes/pages.py
        elif path == "/api/documents":
            if not self._check_admin_auth():
                self._send_json({"error": "Unauthorized"}, status_code=401)
                return
            docs_dir = os.path.join(DATA_DIR, "generated_docs")
            os.makedirs(docs_dir, exist_ok=True)
            documents = []
            for fname in sorted(os.listdir(docs_dir), reverse=True):
                if fname.endswith(".zip"):
                    fpath = os.path.join(docs_dir, fname)
                    try:
                        stat = os.stat(fpath)
                        documents.append(
                            {
                                "filename": fname,
                                "size_bytes": stat.st_size,
                                "created": datetime.datetime.fromtimestamp(
                                    stat.st_mtime
                                ).isoformat(),
                                "download_url": f"/api/documents/{fname}",
                            }
                        )
                    except OSError:
                        continue
            response = json.dumps(
                {"total": len(documents), "documents": documents[:100]}
            )
            self.send_response(200)
            self.send_header("Content-Type", "application/json")
            self.end_headers()
            self.wfile.write(response.encode())
        elif path.startswith("/api/documents/"):
            if not self._check_admin_auth():
                self._send_json(
                    {"error": "Unauthorized - admin key required"}, status_code=401
                )
                return
            fname = path.split("/")[-1]
            # Security: sanitize filename to prevent path traversal
            fname = re.sub(r"[^\w\.\-]", "", fname)
            if not fname or ".." in fname:
                self.send_error(400, "Invalid filename")
                return
            doc_path = os.path.join(DATA_DIR, "generated_docs", fname)
            # Verify the resolved path is within the docs directory
            docs_dir = os.path.realpath(os.path.join(DATA_DIR, "generated_docs"))
            real_path = os.path.realpath(doc_path)
            if not real_path.startswith(docs_dir):
                self.send_error(403, "Access denied")
                return
            if os.path.exists(doc_path) and os.path.isfile(doc_path):
                try:
                    with open(doc_path, "rb") as df:
                        content = df.read()
                    self.send_response(200)
                    self.send_header("Content-Type", "application/zip")
                    self.send_header(
                        "Content-Disposition", f'attachment; filename="{fname}"'
                    )
                    self.send_header("Content-Length", str(len(content)))
                    self.end_headers()
                    self.wfile.write(content)
                except Exception as e:
                    self.send_response(500)
                    self.send_header("Content-Type", "application/json")
                    self.end_headers()
                    logger.error("Document read error: %s", e)
                    self.wfile.write(
                        _error_response("Failed to read document", "SERVER_ERROR", 500)[
                            0
                        ]
                    )
            else:
                self.send_response(404)
                self.send_header("Content-Type", "application/json")
                self.end_headers()
                self.wfile.write(
                    _error_response("Document not found", "NOT_FOUND", 404)[0]
                )
        # NOTE: /nova-jarvis, /nova now handled by routes/pages.py
        elif path.startswith("/nova/shared/"):
            # ── Shared conversation view (read-only) ──
            share_id = path.split("/nova/shared/")[-1].rstrip("/")
            if not share_id:
                self.send_error(404, "Share ID required")
                return
            # Load from Supabase
            shared_data: dict | None = None
            try:
                result = _supabase_rest(
                    "nova_shared_conversations",
                    params=f"?share_id=eq.{urllib.parse.quote(share_id)}&limit=1",
                )
                if result and len(result) > 0:
                    shared_data = result[0]
            except Exception as _sle:
                logger.error(
                    "Load shared conversation failed (table nova_shared_conversations may not exist): %s",
                    _sle,
                    exc_info=True,
                )
            if not shared_data:
                self.send_error(404, "Shared conversation not found")
                return
            title = shared_data.get("title") or "Shared Conversation"
            messages_json = shared_data.get("messages") or "[]"
            try:
                messages = (
                    json.loads(messages_json)
                    if isinstance(messages_json, str)
                    else messages_json
                )
            except json.JSONDecodeError:
                messages = []
            # Build read-only HTML from nova.html template concept
            msgs_html = ""
            for msg in messages:
                role = msg.get("role") or "user"
                content = msg.get("content") or ""
                is_user = role == "user"
                avatar_cls = "user" if is_user else "nova"
                avatar_txt = "U" if is_user else "N"
                sender_txt = "You" if is_user else "Nova"
                # Escape content for display
                safe_content = (
                    content.replace("&", "&amp;")
                    .replace("<", "&lt;")
                    .replace(">", "&gt;")
                )
                msgs_html += (
                    f'<div class="message" style="animation:none;opacity:1">'
                    f'<div class="message-avatar {avatar_cls}">{avatar_txt}</div>'
                    f'<div class="message-body">'
                    f'<div class="message-sender {avatar_cls}">{sender_txt}</div>'
                    f'<div class="message-content"><p>{safe_content}</p></div>'
                    f"</div></div>"
                )
            escaped_title = (
                title.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            )
            shared_html = f"""<!doctype html>
<html lang="en"><head><meta charset="UTF-8"/><meta name="viewport" content="width=device-width, initial-scale=1.0"/>
<title>Nova AI - {escaped_title}</title>
<link href="https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700&display=swap" rel="stylesheet"/>
<style>
:root {{--bg-primary:#0f0f1a;--bg-secondary:#1a1a2e;--text-primary:#e4e4e7;--text-secondary:#888;
--accent-teal:#6bb3cd;--accent-purple:#5a54bd;--accent-purple-light:#7b75d4;--border-color:rgba(107,179,205,0.1);}}
*,*::before,*::after {{margin:0;padding:0;box-sizing:border-box;}}
body {{background:var(--bg-primary);color:var(--text-primary);font-family:'Inter',sans-serif;font-size:15px;line-height:1.6;}}
.container {{max-width:720px;margin:0 auto;padding:40px 20px;}}
.header {{text-align:center;margin-bottom:32px;padding-bottom:16px;border-bottom:1px solid var(--border-color);}}
.header h1 {{font-size:20px;color:#fff;margin-bottom:4px;}}
.header p {{font-size:13px;color:var(--text-secondary);}}
.message {{display:flex;gap:12px;margin-bottom:24px;}}
.message-avatar {{width:32px;height:32px;border-radius:50%;flex-shrink:0;display:flex;align-items:center;justify-content:center;font-size:13px;font-weight:600;}}
.message-avatar.nova {{background:linear-gradient(135deg,var(--accent-teal),#4a9db5);color:#fff;}}
.message-avatar.user {{background:var(--accent-purple);color:#fff;}}
.message-body {{flex:1;min-width:0;}}
.message-sender {{font-size:13px;font-weight:600;margin-bottom:4px;}}
.message-sender.nova {{color:var(--accent-teal);}}
.message-sender.user {{color:var(--accent-purple-light);}}
.message-content {{font-size:15px;line-height:1.7;word-wrap:break-word;}}
.message-content p {{margin:0 0 12px;}} .message-content p:last-child {{margin-bottom:0;}}
.back-link {{display:inline-block;margin-top:24px;padding:8px 16px;border-radius:8px;border:1px solid var(--border-color);color:var(--accent-teal);font-size:13px;text-decoration:none;transition:all 0.15s;}}
.back-link:hover {{background:rgba(107,179,205,0.08);text-decoration:none;}}
</style></head><body>
<div class="container">
<div class="header"><h1>{escaped_title}</h1><p>Shared Nova AI conversation (read-only)</p></div>
{msgs_html}
<a href="/nova" class="back-link">Open Nova AI</a>
</div></body></html>"""
            self.send_response(200)
            self.send_header("Content-Type", "text/html; charset=utf-8")
            self.end_headers()
            self.wfile.write(shared_html.encode())
        # NOTE: /admin/nova now handled by routes/pages.py
        # ── Embedded ATS Widget routes ──
        elif path == "/embed/nova.js":
            # Serve the embed widget JS with permissive CORS for cross-origin embedding
            js_path = os.path.join(
                os.path.dirname(os.path.abspath(__file__)), "static", "nova-embed.js"
            )
            try:
                with open(js_path, "rb") as f:
                    js_content = f.read()
                with _embed_stats_lock:
                    _embed_stats["js_serves"] += 1
                self.send_response(200)
                self.send_header(
                    "Content-Type", "application/javascript; charset=utf-8"
                )
                self.send_header("Content-Length", str(len(js_content)))
                self.send_header("Access-Control-Allow-Origin", "*")
                self.send_header("Access-Control-Allow-Methods", "GET")
                self.send_header("Cache-Control", "public, max-age=3600")  # 1 hour
                self.end_headers()
                self.wfile.write(js_content)
            except FileNotFoundError:
                self.send_error(404, "Embed widget not found")
        elif path == "/embed/demo":
            # Serve the embed demo page
            demo_path = os.path.join(
                os.path.dirname(os.path.abspath(__file__)),
                "static",
                "nova-embed-demo.html",
            )
            try:
                with open(demo_path, "rb") as f:
                    demo_content = f.read()
                with _embed_stats_lock:
                    _embed_stats["demo_views"] += 1
                self.send_response(200)
                self.send_header("Content-Type", "text/html; charset=utf-8")
                self.send_header("Content-Length", str(len(demo_content)))
                self.end_headers()
                self.wfile.write(demo_content)
            except FileNotFoundError:
                self.send_error(404, "Demo page not found")
        elif path == "/api/embed/stats":
            # Embed stats beacon endpoint (1x1 pixel or JSON based on Accept header)
            query_params = urllib.parse.parse_qs(parsed.query)
            event = (query_params.get("event", ["load"])[0]) or "load"
            host = (query_params.get("host", [""])[0]) or ""
            _record_embed_event(event, host)
            # Return 1x1 transparent GIF (for pixel tracking from img tags)
            gif_bytes = bytes(
                [
                    0x47,
                    0x49,
                    0x46,
                    0x38,
                    0x39,
                    0x61,  # GIF89a
                    0x01,
                    0x00,
                    0x01,
                    0x00,
                    0x80,
                    0xFF,  # 1x1 with color table
                    0x00,
                    0xFF,
                    0xFF,
                    0xFF,
                    0x00,
                    0x00,  # white + black
                    0x00,
                    0x21,
                    0xF9,
                    0x04,
                    0x01,
                    0x00,  # graphic control ext
                    0x00,
                    0x00,
                    0x00,
                    0x2C,
                    0x00,
                    0x00,  # image descriptor
                    0x00,
                    0x00,
                    0x01,
                    0x00,
                    0x01,
                    0x00,
                    0x00,
                    0x02,
                    0x02,
                    0x44,
                    0x01,
                    0x00,  # image data
                    0x3B,  # trailer
                ]
            )
            self.send_response(200)
            self.send_header("Content-Type", "image/gif")
            self.send_header("Content-Length", str(len(gif_bytes)))
            self.send_header("Access-Control-Allow-Origin", "*")
            self.send_header("Cache-Control", "no-store")
            self.end_headers()
            self.wfile.write(gif_bytes)
        elif path.startswith("/static/"):
            # Serve static files (JS, CSS, images) from the static/ directory
            # URL-decode the path first to catch encoded traversal (%2e%2e)
            decoded_path = urllib.parse.unquote(path)
            safe_path = decoded_path.lstrip("/")
            # Security: reject directory traversal attempts
            if ".." in safe_path or safe_path.startswith("/") or "\x00" in safe_path:
                self.send_error(403)
                return
            filepath = os.path.join(
                os.path.dirname(os.path.abspath(__file__)), safe_path
            )
            # Resolve symlinks and verify the real path stays within static/
            static_root = os.path.realpath(
                os.path.join(os.path.dirname(os.path.abspath(__file__)), "static")
            )
            real_filepath = os.path.realpath(filepath)
            if (
                not real_filepath.startswith(static_root + os.sep)
                and real_filepath != static_root
            ):
                self.send_error(403)
                return
            ext = os.path.splitext(filepath)[1].lower()
            content_types = {
                ".js": "application/javascript",
                ".css": "text/css",
                ".html": "text/html",
                ".json": "application/json",
                ".png": "image/png",
                ".jpg": "image/jpeg",
                ".svg": "image/svg+xml",
                ".ico": "image/x-icon",
            }
            ctype = content_types.get(ext, "application/octet-stream")
            if os.path.isfile(filepath):
                try:
                    # Serve minified version if available (CSS/JS only)
                    url_path = f"/{safe_path}"
                    minified = get_minified(url_path)
                    if minified is not None:
                        content, etag = minified
                    else:
                        with open(filepath, "rb") as f:
                            content = f.read()
                        etag = hashlib.sha256(content).hexdigest()[:12]

                    # Cache policy: JS/CSS use must-revalidate so browsers
                    # always check ETags (304 on match); images/fonts use
                    # long-lived immutable cache (they rarely change).
                    if ext in (".js", ".css"):
                        cache_control = "public, no-cache, must-revalidate"
                    else:
                        cache_control = "public, max-age=31536000, immutable"

                    # ETag-based conditional response (304 Not Modified)
                    client_etag = self.headers.get("If-None-Match", "")
                    if client_etag and client_etag.strip('"') == etag:
                        self.send_response(304)
                        self.send_header("ETag", f'"{etag}"')
                        self.send_header("Cache-Control", cache_control)
                        self.end_headers()
                        return

                    self.send_response(200)
                    self.send_header("Content-Type", f"{ctype}; charset=utf-8")
                    self.send_header("Content-Length", str(len(content)))
                    self.send_header("Cache-Control", cache_control)
                    self.send_header("ETag", f'"{etag}"')
                    self.end_headers()
                    self.wfile.write(content)
                except FileNotFoundError:
                    self.send_error(404)
            else:
                self.send_error(404)
        # ── Feature 1: Swagger UI + OpenAPI spec ──
        elif path == "/docs":
            body = _SWAGGER_UI_HTML.encode("utf-8")
            self.send_response(200)
            self.send_header("Content-Type", "text/html; charset=utf-8")
            self.send_header("Content-Length", str(len(body)))
            self.send_header("Cache-Control", "public, max-age=3600")
            self.end_headers()
            self.wfile.write(body)
        elif path == "/api/docs/openapi.json":
            body = json.dumps(_OPENAPI_SPEC, indent=2).encode("utf-8")
            self.send_response(200)
            self.send_header("Content-Type", "application/json")
            self.send_header("Content-Length", str(len(body)))
            self.send_header("Cache-Control", "public, max-age=3600")
            self.end_headers()
            self.wfile.write(body)
        # ── Feature 2c: Async job polling ──
        elif path.startswith("/api/jobs/"):
            job_id = path.split("/")[-1]
            if not job_id or not re.match(r"^[a-f0-9]{1,12}$", job_id):
                self.send_response(400)
                self.send_header("Content-Type", "application/json")
                self.end_headers()
                self.wfile.write(
                    _error_response("Invalid job_id", "VALIDATION_ERROR", 400)[0]
                )
                return
            now = time.time()
            with _generation_jobs_lock:
                job = _generation_jobs.get(job_id)
                if not job:
                    # S47: Fallback to Supabase for restart-safe Slack download links
                    try:
                        from supabase_client import get_client as _sb_dl_client
                        import base64 as _b64_dl

                        _sb_dl = _sb_dl_client()
                        if _sb_dl:
                            _sb_result = (
                                _sb_dl.table("nova_generated_plans")
                                .select("*")
                                .eq("job_id", job_id)
                                .execute()
                            )
                            if _sb_result.data:
                                _row = _sb_result.data[0]
                                _zip_bytes = _b64_dl.b64decode(_row["zip_data"])
                                _dl_filename = _row.get("filename") or "Media_Plan.zip"
                                _dl_ct = "application/zip"
                                accept_hdr = self.headers.get("Accept") or ""
                                if "application/json" in accept_hdr:
                                    # Poll mode: tell frontend it's ready
                                    self._send_json(
                                        {
                                            "job_id": job_id,
                                            "status": "completed",
                                            "progress_pct": 100,
                                            "status_message": "Complete",
                                            "filename": _dl_filename,
                                            "content_type": _dl_ct,
                                            "download_url": f"/api/jobs/{job_id}",
                                            "source": "supabase",
                                        }
                                    )
                                    return
                                self.send_response(200)
                                self.send_header("Content-Type", _dl_ct)
                                self.send_header(
                                    "Content-Disposition",
                                    f'attachment; filename="{_dl_filename}"',
                                )
                                self.send_header("Content-Length", str(len(_zip_bytes)))
                                self.end_headers()
                                self.wfile.write(_zip_bytes)
                                logger.info(
                                    "S47: Served plan from Supabase for job %s", job_id
                                )
                                return
                    except Exception as _sb_dl_err:
                        logger.debug(
                            "S47: Supabase fallback lookup failed: %s", _sb_dl_err
                        )
                    self.send_response(404)
                    self.send_header("Content-Type", "application/json")
                    self.end_headers()
                    self.wfile.write(
                        _error_response("Job not found or expired", "NOT_FOUND", 404)[0]
                    )
                    return
                # Bug #14 fix: Validate session owns this job (IDOR prevention)
                # C1 fix: Use nova_session (HttpOnly) cookie, fallback to csrf_token
                _job_session = job.get("_session_token") or ""
                if _job_session:
                    _poll_csrf = _parse_cookie_value(
                        self.headers.get("Cookie") or "", "nova_session"
                    ) or _parse_cookie_value(
                        self.headers.get("Cookie") or "", "csrf_token"
                    )
                    if not hmac.compare_digest(_job_session, _poll_csrf):
                        self.send_response(403)
                        self.send_header("Content-Type", "application/json")
                        self.end_headers()
                        self.wfile.write(
                            _error_response(
                                "Access denied: job belongs to a different session",
                                "FORBIDDEN",
                                403,
                            )[0]
                        )
                        return
                # Auto-expire check
                if now - job["created"] > _GENERATION_JOB_EXPIRY_SECONDS:
                    _generation_jobs.pop(job_id, None)
                    self.send_response(404)
                    self.send_header("Content-Type", "application/json")
                    self.end_headers()
                    self.wfile.write(
                        _error_response("Job expired", "NOT_FOUND", 410)[0]
                    )
                    return
                job_status = job["status"]
                elapsed = round(now - job["created"], 1)
            if job_status == "completed":
                accept_header = self.headers.get("Accept") or ""
                if "application/json" in accept_header:
                    # Poll mode: return JSON status so frontend knows it's ready
                    filename = job.get("result_filename") or "result.zip"
                    content_type = job.get("result_content_type") or "application/zip"
                    self._send_json(
                        {
                            "job_id": job_id,
                            "status": "completed",
                            "progress_pct": 100,
                            "status_message": "Complete",
                            "filename": filename,
                            "content_type": content_type,
                            "download_url": f"/api/jobs/{job_id}",
                            "elapsed_seconds": elapsed,
                        }
                    )
                else:
                    # Download mode: return the binary content
                    result_bytes = job.get("result_bytes", b"")
                    content_type = job.get("result_content_type", "application/zip")
                    filename = job.get("result_filename", "result.zip")
                    # S37 fix: If result_bytes already cleared (retry after
                    # previous download), return 410 so frontend can show a
                    # helpful message instead of silently serving 0 bytes.
                    if not result_bytes:
                        self._send_json(
                            {
                                "job_id": job_id,
                                "status": "expired",
                                "error": "File already downloaded. Please regenerate.",
                            },
                            status_code=410,
                        )
                        return
                    self.send_response(200)
                    self.send_header("Content-Type", content_type)
                    self.send_header(
                        "Content-Disposition", f'attachment; filename="{filename}"'
                    )
                    self.send_header("Content-Length", str(len(result_bytes)))
                    self.end_headers()
                    self.wfile.write(result_bytes)
                    # S37 fix: Keep result_bytes for 5 min after first download
                    # to allow retries if network drops mid-transfer. Previously
                    # bytes were cleared immediately, causing "failed download"
                    # on retry (Ashlie's bug). Deferred cleanup via timer thread.
                    with _generation_jobs_lock:
                        if job_id in _generation_jobs:
                            _generation_jobs[job_id]["_downloaded"] = True
                            _generation_jobs[job_id]["_download_time"] = time.time()

                            # Schedule deferred cleanup: free bytes after 5 min
                            def _deferred_bytes_cleanup(_jid: str) -> None:
                                time.sleep(300)  # 5 minutes
                                with _generation_jobs_lock:
                                    j = _generation_jobs.get(_jid)
                                    if j and j.get("_downloaded"):
                                        j["result_bytes"] = b""

                            t = threading.Thread(
                                target=_deferred_bytes_cleanup,
                                args=(job_id,),
                                daemon=True,
                            )
                            t.start()
            elif job_status == "failed":
                err_msg = job.get("error", "Generation failed")
                self._send_json(
                    {
                        "job_id": job_id,
                        "status": "failed",
                        "error": err_msg,
                        "created": datetime.datetime.fromtimestamp(
                            job["created"]
                        ).isoformat(),
                        "elapsed_seconds": elapsed,
                    }
                )
            else:
                self._send_json(
                    {
                        "job_id": job_id,
                        "status": "processing",
                        "progress_pct": job.get("progress_pct") or 0,
                        "status_message": job.get("status_message") or "Processing...",
                        "created": datetime.datetime.fromtimestamp(
                            job["created"]
                        ).isoformat(),
                        "elapsed_seconds": elapsed,
                    }
                )
        # ── Plan Results JSON (on-screen dashboard data) ──
        elif path.startswith("/api/plan-results/"):
            _pr_id = path.split("/")[-1]
            if not _pr_id or not re.match(r"^[a-f0-9]{1,12}$", _pr_id):
                self._send_json({"error": "Invalid plan_id"}, status_code=400)
                return
            now = time.time()
            with _plan_results_lock:
                _pr_entry = _plan_results_store.get(_pr_id)
                if not _pr_entry:
                    self._send_json(
                        {"error": "Plan not found or expired"}, status_code=404
                    )
                    return
                if now - _pr_entry["created"] > _PLAN_RESULTS_TTL_SECONDS:
                    _plan_results_store.pop(_pr_id, None)
                    self._send_json({"error": "Plan expired"}, status_code=404)
                    return
                self._send_json({"plan_id": _pr_id, **_pr_entry["data"]})
        # ── Google Sheets URL for a generated plan ──
        elif path == "/api/plan/sheets-url":
            _qs_parsed = urllib.parse.parse_qs(urllib.parse.urlparse(self.path).query)
            _su_plan_id = (_qs_parsed.get("plan_id") or [""])[0]
            if not _su_plan_id or not re.match(r"^[a-f0-9]{1,12}$", _su_plan_id):
                self._send_json(
                    {"error": "Missing or invalid plan_id"}, status_code=400
                )
                return
            with _plan_results_lock:
                _su_entry = _plan_results_store.get(_su_plan_id)
                if not _su_entry:
                    self._send_json(
                        {"error": "Plan not found or expired"}, status_code=404
                    )
                    return
                _su_url = _su_entry.get("sheets_url") or ""
            if _su_url:
                self._send_json(
                    {"plan_id": _su_plan_id, "sheets_url": _su_url, "ready": True}
                )
            else:
                # Sheet may still be generating in the background
                self._send_json(
                    {"plan_id": _su_plan_id, "sheets_url": None, "ready": False}
                )
        # ── Plan events list: GET /api/plan/events (no trailing plan_id) ──
        elif path == "/api/plan/events":
            try:
                from plan_events import get_event_store as _get_pe_list_store

                _pe_list = _get_pe_list_store()
                _pe_stats = _pe_list.get_stats()
                # Include recent events per plan (last 10 each)
                _pe_recent: dict[str, list[dict]] = {}
                for _pe_pid in list(_pe_stats.get("plans", {}).keys())[:50]:
                    _pe_evts = _pe_list.get_events(_pe_pid)
                    _pe_recent[_pe_pid] = [e.to_dict() for e in _pe_evts[-10:]]
                self._send_json(
                    {
                        "total_plans": _pe_stats.get("total_plans", 0),
                        "total_events": _pe_stats.get("total_events", 0),
                        "created_at": _pe_stats.get("created_at") or "",
                        "plans_summary": _pe_stats.get("plans", {}),
                        "recent_events": _pe_recent,
                    }
                )
            except ImportError:
                self._send_json(
                    {"error": "plan_events module not available"}, status_code=503
                )
            except Exception as e:
                logger.error(f"Plan events list error: {e}", exc_info=True)
                self._send_json({"error": _internal_error_msg()}, status_code=500)
        # ── Event-sourced plan state (audit trail, time-travel) ──
        elif path.startswith("/api/plan/events/"):
            _es_plan_id = path.split("/")[-1]
            if not _es_plan_id:
                self._send_json({"error": "Missing plan_id"}, status_code=400)
                return
            try:
                from event_store import get_event_store

                _es = get_event_store()
                events = _es.get_events(_es_plan_id)
                self._send_json(
                    {
                        "plan_id": _es_plan_id,
                        "event_count": len(events),
                        "events": [e.to_dict() for e in events],
                    }
                )
            except ImportError:
                self._send_json(
                    {"error": "event_store module not available"}, status_code=503
                )
            except Exception as e:
                logger.error(f"Event store GET error: {e}", exc_info=True)
                self._send_json({"error": _internal_error_msg()}, status_code=500)
        elif path.startswith("/api/plan/state/"):
            _es_plan_id = path.split("/")[-1]
            if not _es_plan_id:
                self._send_json({"error": "Missing plan_id"}, status_code=400)
                return
            try:
                from event_store import get_event_store

                _es = get_event_store()
                _qs = urllib.parse.parse_qs(parsed.query)
                _ver_str = _qs.get("version", [None])[0]
                if _ver_str is not None:
                    try:
                        _ver = int(_ver_str)
                    except (ValueError, TypeError):
                        self._send_json(
                            {"error": "version must be an integer"}, status_code=400
                        )
                        return
                    state = _es.get_state_at(_es_plan_id, _ver)
                else:
                    state = _es.get_current_state(_es_plan_id)
                status = 404 if "error" in state else 200
                self._send_json(state, status_code=status)
            except ImportError:
                self._send_json(
                    {"error": "event_store module not available"}, status_code=503
                )
            except Exception as e:
                logger.error(f"Event store state GET error: {e}", exc_info=True)
                self._send_json({"error": _internal_error_msg()}, status_code=500)
        # ── Event-sourced plan state machine (plan_events module) ──
        elif path.startswith("/api/plans/") and path.endswith("/history"):
            # GET /api/plans/<plan_id>/history
            parts = path.split("/")
            # Expected: ['', 'api', 'plans', '<plan_id>', 'history']
            _pe_plan_id = parts[3] if len(parts) >= 5 else ""
            if not _pe_plan_id:
                self._send_json({"error": "Missing plan_id"}, status_code=400)
                return
            try:
                from plan_events import get_event_store as _get_pe_store

                _pe_store = _get_pe_store()
                _pe_qs = urllib.parse.parse_qs(parsed.query)
                _pe_since = _pe_qs.get("since_version", [None])[0]
                since_v = int(_pe_since) if _pe_since else 0
                events = _pe_store.get_events(_pe_plan_id, since_version=since_v)
                snapshot = _pe_store.get_snapshot(_pe_plan_id)
                self._send_json(
                    {
                        "plan_id": _pe_plan_id,
                        "event_count": len(events),
                        "events": [e.to_dict() for e in events],
                        "current_snapshot": snapshot,
                    }
                )
            except ImportError:
                self._send_json(
                    {"error": "plan_events module not available"}, status_code=503
                )
            except Exception as e:
                logger.error(f"Plan events history error: {e}", exc_info=True)
                self._send_json({"error": _internal_error_msg()}, status_code=500)
        # NOTE: /plan/{id} shareable view now handled by routes/campaign.py
        # ── Feature 5a: SLO compliance ──
        # ── Platform Observability (aggregated health dashboard) ──
        # ── Feature 5b: Eval scores ──
        # ── Feature 3c: Per-key usage dashboard ──
        # ── Admin Statistics Endpoint ──
        # ── PostHog Analytics Admin Endpoint ──
        # ── LLM Cost Tracking Report ──
        # ── Audit Events Endpoint ──
        # NOTE: /tracker, /simulator, /vendor-iq pages now handled by routes/pages.py

        # ── Saved Plans (Supabase persistence -- Cindy request S47) ──
        elif path == "/api/saved-plans":
            try:
                if not _check_joveo_auth(self):
                    self._send_error("Authentication required", 401)
                    return
                try:
                    from supabase_client import get_client

                    sb = get_client()
                except Exception:
                    sb = None
                if not sb:
                    self._send_json({"plans": []})
                    return
                email = _parse_cookie_value(
                    self.headers.get("Cookie") or "", "nova_user_email"
                )
                result = (
                    sb.table("nova_saved_plans")
                    .select(
                        "id,plan_name,industry,location,budget,created_at,updated_at"
                    )
                    .eq("user_email", email or "unknown")
                    .order("created_at", desc=True)
                    .limit(20)
                    .execute()
                )
                self._send_json({"plans": result.data or []})
            except Exception as exc:
                logger.error("Saved plans list error: %s", exc, exc_info=True)
                self._send_json({"plans": []})

        elif path.startswith("/api/saved-plans/") and not path.endswith("/"):
            # GET /api/saved-plans/<id> -- fetch full plan data
            try:
                if not _check_joveo_auth(self):
                    self._send_error("Authentication required", 401)
                    return
                plan_id = path.split("/")[-1]
                if not plan_id.isdigit():
                    self._send_json({"error": "Invalid plan ID"}, status_code=400)
                    return
                try:
                    from supabase_client import get_client

                    sb = get_client()
                except Exception:
                    sb = None
                if not sb:
                    self._send_json({"error": "Storage unavailable"}, status_code=503)
                    return
                result = (
                    sb.table("nova_saved_plans")
                    .select("*")
                    .eq("id", int(plan_id))
                    .execute()
                )
                if not result.data:
                    self._send_json({"error": "Plan not found"}, status_code=404)
                    return
                self._send_json(result.data[0])
            except Exception as exc:
                logger.error("Saved plan fetch error: %s", exc, exc_info=True)
                self._send_json({"error": "Failed to fetch plan"}, status_code=500)

        # ── Budget Simulator GET APIs ──
        elif path == "/api/simulator/defaults":
            try:
                qs = urllib.parse.parse_qs(parsed.query)
                total_budget = float(qs.get("budget", ["100000"])[0])
                industry = qs.get("industry", ["Technology"])[0]
                roles = qs.get("roles", ["Software Engineer"])[0]
                locations = qs.get("locations", ["United States"])[0]
                from budget_simulator import get_default_allocation

                result = get_default_allocation(
                    total_budget, industry, roles, locations
                )
                self._send_json(result)
            except Exception as e:
                logger.error("Simulator defaults error: %s", e, exc_info=True)
                self._send_json({"error": _internal_error_msg()})
        # NOTE: /competitive through /terms pages now handled by routes/pages.py
        # ── API Integrations Status ──
        elif path == "/api/integrations/status":
            if not _api_integrations_available:
                self._send_json(
                    {
                        "available": False,
                        "error": "api_integrations module not loaded",
                        "apis": {},
                    }
                )
            else:
                try:
                    _api_results = _api_test_all()
                    _passed = sum(1 for v in _api_results.values() if v)
                    _total = len(_api_results)
                    self._send_json(
                        {
                            "available": True,
                            "summary": f"{_passed}/{_total} APIs healthy",
                            "apis": _api_results,
                            "timestamp": datetime.datetime.utcnow().isoformat() + "Z",
                        }
                    )
                except Exception as _e:
                    logger.error(
                        "API integrations status check failed: %s", _e, exc_info=True
                    )
                    self._send_json(
                        {
                            "available": True,
                            "error": "Status check failed",
                            "apis": {},
                        }
                    )
        # ── API Integrations Market Snapshot ──
        elif path == "/api/integrations/market-snapshot":
            qs = urllib.parse.parse_qs(parsed.query)
            _snap_role = qs.get("role", [""])[0]
            _snap_location = qs.get("location", [""])[0]
            if not _snap_role:
                self._send_json({"error": "Missing required 'role' query parameter"})
            elif not _api_integrations_available:
                self._send_json(
                    {
                        "error": "api_integrations module not loaded",
                        "role": _snap_role,
                        "location": _snap_location,
                    }
                )
            else:
                _snapshot: dict = {
                    "role": _snap_role,
                    "location": _snap_location,
                    "timestamp": datetime.datetime.utcnow().isoformat() + "Z",
                }
                # Adzuna job count + salary
                try:
                    if _api_adzuna:
                        _s_count = _api_adzuna.get_job_count(_snap_role, "us")
                        if _s_count:
                            _snapshot["adzuna_job_count"] = _s_count
                        _s_hist = _api_adzuna.get_salary_histogram(_snap_role, "us")
                        if _s_hist:
                            _snapshot["adzuna_salary_histogram"] = _s_hist
                        _s_companies = _api_adzuna.get_top_companies(_snap_role, "us")
                        if _s_companies:
                            _snapshot["adzuna_top_companies"] = _s_companies
                except (urllib.error.URLError, OSError, ValueError, TypeError) as _ae:
                    logger.error(
                        "Adzuna data for market-snapshot failed: %s", _ae, exc_info=True
                    )
                # FRED economic data
                try:
                    if _api_fred:
                        _s_unemp = _api_fred.get_unemployment_rate()
                        if _s_unemp:
                            _snapshot["fred_unemployment"] = _s_unemp
                        _s_cpi = _api_fred.get_cpi_data(months=6)
                        if _s_cpi:
                            _snapshot["fred_cpi"] = _s_cpi
                except (urllib.error.URLError, OSError, ValueError, TypeError) as _fe:
                    logger.error(
                        "FRED data for market-snapshot failed: %s", _fe, exc_info=True
                    )
                # O*NET occupation data
                try:
                    if _api_onet:
                        _s_occs = _api_onet.search_occupations(_snap_role)
                        if _s_occs:
                            _snapshot["onet_occupations"] = _s_occs[:5]
                            _top_code = _s_occs[0].get("code") or ""
                            if _top_code:
                                _s_skills = _api_onet.get_skills(_top_code)
                                if _s_skills:
                                    _snapshot["onet_skills"] = _s_skills[:10]
                except (urllib.error.URLError, OSError, ValueError, TypeError) as _oe:
                    logger.error(
                        "O*NET data for market-snapshot failed: %s", _oe, exc_info=True
                    )
                # BLS employment data
                try:
                    if _api_bls and _api_onet:
                        _s_occs_bls = _snapshot.get("onet_occupations") or []
                        if _s_occs_bls:
                            _s_soc = _s_occs_bls[0].get("code") or ""
                            if _s_soc:
                                _s_oes = _api_bls.get_occupational_employment(_s_soc)
                                if _s_oes:
                                    _snapshot["bls_occupational_employment"] = _s_oes
                except (urllib.error.URLError, OSError, ValueError, TypeError) as _be:
                    logger.error(
                        "BLS data for market-snapshot failed: %s", _be, exc_info=True
                    )
                # Census population
                try:
                    if _api_census:
                        _s_pop = _api_census.get_population_by_state()
                        if _s_pop:
                            _snapshot["census_population"] = _s_pop
                except (urllib.error.URLError, OSError, ValueError, TypeError) as _ce:
                    logger.error(
                        "Census data for market-snapshot failed: %s", _ce, exc_info=True
                    )
                self._send_json(_snapshot)
        # ── API Portal GET routes ──
        elif path.startswith("/api/portal/"):
            if not self._check_admin_auth():
                self._send_json({"error": "Unauthorized"}, status_code=401)
                return
            try:
                from api_portal import handle_portal_api

                content_len = int(self.headers.get("Content-Length") or 0)
                body = (
                    json.loads(self.rfile.read(content_len)) if content_len > 0 else {}
                )
                headers = {k: v for k, v in self.headers.items()}
                status_code, result = handle_portal_api(
                    path.replace("/api/", ""), "GET", body, headers
                )
                self.send_response(status_code)
                self.send_header("Content-Type", "application/json")
                self.end_headers()
                self.wfile.write(json.dumps(result, default=str).encode())
            except Exception as e:
                logger.error("API Portal GET error: %s", e, exc_info=True)
                self._send_json({"error": _internal_error_msg()})
        # ── Market Pulse GET APIs ──
        elif (
            path.startswith("/api/pulse/")
            and not path.startswith("/api/pulse/generate")
            and not path.startswith("/api/pulse/send")
            and not path.startswith("/api/pulse/scheduler")
        ):
            try:
                from market_pulse import handle_pulse_api

                status_code, result = handle_pulse_api(path, method="GET")
                if status_code == 200:
                    self._send_json(result)
                else:
                    self.send_response(status_code)
                    self.send_header("Content-Type", "application/json")
                    self.end_headers()
                    self.wfile.write(json.dumps(result).encode())
            except Exception as e:
                logger.error("Market Pulse GET error: %s", e, exc_info=True)
                self._send_json({"ok": False, "error": _internal_error_msg()})
        # ── Firecrawl Status (S72: removed; endpoint kept for backwards compat) ──
        elif path == "/api/firecrawl/status":
            self._send_json(
                {
                    "configured": False,
                    "removed": True,
                    "removed_in": "S72",
                    "reason": "Firecrawl module deleted (credits exhausted). "
                    "Use /api/scraper/status for active scraper tiers.",
                }
            )
        # ── Web Scraper Router Status (multi-tier) ──
        elif path == "/api/scraper/status":
            try:
                from web_scraper_router import get_scraper_status

                self._send_json(get_scraper_status())
            except ImportError:
                self._send_json(
                    {
                        "error": "web_scraper_router module not available",
                        "total_tiers": 0,
                    }
                )
            except Exception as exc:
                logger.error("Scraper status error: %s", exc, exc_info=True)
                self._send_json({"error": _internal_error_msg()}, status_code=500)
        # ── Sentry Integration Status + Recent Issues ──
        elif path == "/api/sentry/issues":
            if not self._check_admin_auth():
                self._send_json({"error": "Unauthorized"}, status_code=401)
                return
            try:
                if _sentry_integration_available:
                    status = _get_sentry_status()
                    # Optionally fetch live issues from Sentry API
                    qs = urllib.parse.parse_qs(parsed.query)
                    if qs.get("live", ["false"])[0].lower() == "true":
                        status["live_issues"] = _SentryAPIClient.fetch_recent_issues(
                            hours=int(qs.get("hours", ["24"])[0])
                        )
                    self._send_json(status)
                else:
                    self._send_json(
                        {
                            "configured": False,
                            "error": "sentry_integration module not available",
                        }
                    )
            except Exception as e:
                logger.error("Sentry issues endpoint error: %s", e, exc_info=True)
                self._send_json({"error": _internal_error_msg()}, status_code=500)
        # ── Market Pulse News (S72: Firecrawl primary removed; Tavily-only) ──
        elif path == "/api/market-pulse/news":
            try:
                qs = urllib.parse.parse_qs(parsed.query)
                topic = qs.get("topic", ["recruitment advertising"])[0]
                # S72: fetch_recruitment_news (Firecrawl) deleted. Start with
                # an empty articles list and let Tavily / FRED / BEA / Supabase
                # populate the response in the enrichment block below.
                articles: list = []
                _mp_news_resp: dict = {
                    "ok": True,
                    "articles": articles,
                    "count": 0,
                    "primary_source": "disabled",
                    "note": "Firecrawl news removed in S72; using Tavily + econ data.",
                }
                if True:  # preserve indentation of downstream enrichment block
                    # ── API Integrations: Market Pulse News enrichment ──
                    if _api_integrations_available:
                        _mp_econ: dict = {}
                        try:
                            if _api_fred:
                                _fred_unemp = _api_fred.get_unemployment_rate()
                                if _fred_unemp:
                                    _mp_econ["fred_unemployment"] = _fred_unemp
                                _fred_gdp = _api_fred.get_gdp_growth()
                                if _fred_gdp:
                                    _mp_econ["fred_gdp_growth"] = _fred_gdp
                                if _mp_econ:
                                    logger.info(
                                        "Enriched /api/market-pulse/news with fred economic data"
                                    )
                        except (
                            urllib.error.URLError,
                            OSError,
                            ValueError,
                            TypeError,
                        ) as _fe:
                            logger.error(
                                "FRED enrichment for market-pulse/news failed: %s",
                                _fe,
                                exc_info=True,
                            )
                        try:
                            if _api_bea:
                                _bea_gdp = _api_bea.get_gdp_by_state()
                                if _bea_gdp:
                                    _mp_econ["bea_gdp_by_state"] = _bea_gdp
                                    logger.info(
                                        "Enriched /api/market-pulse/news with bea GDP data"
                                    )
                        except (
                            urllib.error.URLError,
                            OSError,
                            ValueError,
                            TypeError,
                        ) as _be:
                            logger.error(
                                "BEA enrichment for market-pulse/news failed: %s",
                                _be,
                                exc_info=True,
                            )
                        if _mp_econ:
                            _mp_news_resp["economic_context"] = _mp_econ
                    # ── Tavily: Real-time recruitment industry news ──
                    if _tavily_available and _tavily_recruitment_news:
                        try:
                            _tav_news = _tavily_recruitment_news(topic)
                            if _tav_news:
                                _mp_news_resp["tavily_news"] = _tav_news
                                logger.info(
                                    "Enriched /api/market-pulse/news with tavily news (%d articles)",
                                    len(_tav_news),
                                )
                        except (
                            urllib.error.URLError,
                            OSError,
                            ValueError,
                            TypeError,
                        ) as _te:
                            logger.error(
                                "Tavily enrichment for market-pulse/news failed: %s",
                                _te,
                                exc_info=True,
                            )

                    # Enrich with Supabase market trends
                    if _supabase_data_available:
                        try:
                            _mp_news_trends = get_market_trends(category="cpc_trends")
                            if _mp_news_trends:
                                _mp_news_resp["supabase_market_trends"] = (
                                    _mp_news_trends
                                )
                        except (urllib.error.URLError, OSError, ValueError) as sb_err:
                            logger.error(
                                f"Supabase market trends for news failed: {sb_err}",
                                exc_info=True,
                            )
                    self._send_json(_mp_news_resp)
            except Exception as e:
                logger.error("Market pulse news error: %s", e, exc_info=True)
                self._send_json(
                    {"ok": False, "error": _internal_error_msg(), "articles": []}
                )
        # ── ElevenLabs Health Check ──
        elif path == "/api/elevenlabs/health":
            if not _elevenlabs_available:
                self._send_json(
                    {
                        "healthy": False,
                        "error": "elevenlabs_integration module not available",
                    }
                )
            else:
                try:
                    health = _elevenlabs_health()
                    self._send_json(health)
                except Exception as e:
                    logger.error("ElevenLabs health check error: %s", e, exc_info=True)
                    self._send_json(
                        {
                            "healthy": False,
                            "error": f"Health check failed: {e}",
                        }
                    )
        # ── Proactive Intelligence Insights (GET) ──
        elif path == "/api/insights":
            try:
                from nova_proactive import get_insights, get_unread_insights

                params = urllib.parse.parse_qs(urllib.parse.urlparse(self.path).query)
                if "unread" in params:
                    self._send_json({"insights": get_unread_insights()})
                else:
                    limit = int(params.get("limit", ["20"])[0])
                    self._send_json({"insights": get_insights(limit)})
            except Exception as e:
                logger.error("Insights list error: %s", e, exc_info=True)
                self._send_json({"insights": [], "error": _internal_error_msg()})

        elif path == "/api/insights/stats":
            try:
                from nova_proactive import get_proactive_stats

                self._send_json(get_proactive_stats())
            except Exception as e:
                logger.error("Proactive stats error: %s", e, exc_info=True)
                self._send_json({"error": _internal_error_msg()}, status_code=500)

        # ── Outcome Pipeline: prediction accuracy report ──
        elif path == "/api/outcomes/accuracy":
            try:
                from outcome_pipeline import get_accuracy_report

                params = urllib.parse.parse_qs(parsed.query)
                role_family = (params.get("role_family", [""])[0] or "").strip()
                time_range = int(params.get("time_range_days", ["90"])[0])
                report = get_accuracy_report(
                    role_family=role_family, time_range_days=time_range
                )
                self._send_json(report)
            except ValueError as ve:
                self._send_json({"error": f"Invalid parameter: {ve}"}, status_code=400)
            except Exception as e:
                logger.error("Outcome accuracy report error: %s", e, exc_info=True)
                self._send_json(
                    {"error": f"Accuracy report failed: {e}"}, status_code=500
                )

        # ── Outcome Pipeline: get outcomes for a plan ──
        elif path.startswith("/api/outcomes/") and not path.endswith("/record"):
            try:
                from outcome_pipeline import (
                    get_plan_outcomes as _get_plan_outcomes,
                    generate_improvement_suggestions as _gen_suggestions,
                )

                plan_id = path[len("/api/outcomes/") :]
                if not plan_id:
                    self._send_json(
                        {"error": "Missing plan_id in URL"}, status_code=400
                    )
                else:
                    funnel = _get_plan_outcomes(plan_id)
                    suggestions = _gen_suggestions(plan_id)
                    funnel["improvement_suggestions"] = suggestions
                    self._send_json(funnel)
            except Exception as e:
                logger.error("Outcome lookup error: %s", e, exc_info=True)
                self._send_json(
                    {"error": f"Outcome lookup failed: {e}"}, status_code=500
                )

        # ── Role Taxonomy (semantic similarity) ──
        elif path == "/api/roles/similar":
            try:
                from role_taxonomy import get_role_taxonomy

                params = urllib.parse.parse_qs(parsed.query)
                title = (params.get("title", [""])[0] or "").strip()
                if not title:
                    self._send_json(
                        {"error": "Missing required query parameter: title"},
                        status_code=400,
                    )
                else:
                    top_k = min(int(params.get("top_k", ["5"])[0]), 20)
                    taxonomy = get_role_taxonomy()
                    similar = taxonomy.find_similar_roles(title, top_k=top_k)
                    cluster = taxonomy.get_role_cluster(title)
                    self._send_json(
                        {
                            "query": title,
                            "similar_roles": [
                                {"role": r, "similarity": s} for r, s in similar
                            ],
                            "cluster": cluster,
                        }
                    )
            except ValueError as ve:
                self._send_json({"error": f"Invalid parameter: {ve}"}, status_code=400)
            except Exception as e:
                logger.error("Role taxonomy error: %s", e, exc_info=True)
                self._send_json(
                    {"error": f"Role taxonomy lookup failed: {e}"},
                    status_code=500,
                )

        # ── Edge Router: Suggest optimal provider ──
        elif path == "/api/routing/suggest":
            try:
                from edge_router import get_optimal_route

                params = urllib.parse.parse_qs(parsed.query)
                query_text = (params.get("query", [""])[0] or "").strip()
                if not query_text:
                    self._send_json(
                        {"error": "Missing required query parameter: query"},
                        status_code=400,
                    )
                else:
                    user_tier = (params.get("tier", ["free"])[0] or "free").strip()
                    route = get_optimal_route(query_text, user_tier=user_tier)
                    route["query_preview"] = (
                        query_text[:80] + "..." if len(query_text) > 80 else query_text
                    )
                    self._send_json(route)
            except ImportError:
                self._send_json(
                    {"error": "edge_router module not available"},
                    status_code=503,
                )
            except Exception as e:
                logger.error("Edge router suggest error: %s", e, exc_info=True)
                self._send_json(
                    {"error": f"Edge router failed: {e}"},
                    status_code=500,
                )

        # ── Anonymized Benchmarking Network ──
        elif path == "/api/benchmarks":
            try:
                from benchmarking import get_benchmarks as _get_bench

                params = urllib.parse.parse_qs(parsed.query)
                role_family = (params.get("role_family", [""])[0] or "").strip() or None
                location = (params.get("location", [""])[0] or "").strip() or None
                self._send_json(_get_bench(role_family, location))
            except Exception as e:
                logger.error("Benchmarking API error: %s", e, exc_info=True)
                self._send_json(
                    {"error": f"Benchmarking lookup failed: {e}", "sample_size": 0},
                    status_code=500,
                )

        # ── CPC/CPA Alerts (GET /api/cpc-alerts) ── S46
        elif path == "/api/cpc-alerts":
            with _cpc_alerts_lock:
                alerts = list(_cpc_alerts)
            self._send_json({"alerts": alerts, "count": len(alerts)})

        # ── Data Freshness Report (GET /api/data/freshness) ──
        elif path == "/api/data/freshness":
            try:
                from kb_loader import get_data_freshness_report

                freshness = get_data_freshness_report()
                self._send_json(freshness)
            except ImportError:
                self._send_json(
                    {"error": "kb_loader module not available"},
                    status_code=503,
                )
            except Exception as e:
                logger.error("Data freshness report error: %s", e, exc_info=True)
                self._send_json(
                    {"error": f"Freshness check failed: {e}"},
                    status_code=500,
                )

        # ── Admin Data Freshness (GET /api/admin/data-freshness) ── S47
        # Same report as /api/data/freshness but under admin auth.
        elif path == "/api/admin/data-freshness":
            if not self._check_admin_auth():
                self._send_error("Unauthorized", "AUTH_REQUIRED", 401)
                return
            try:
                from kb_loader import get_data_freshness_report

                freshness = get_data_freshness_report()
                self._send_json(freshness)
            except ImportError:
                self._send_json(
                    {"error": "kb_loader module not available"},
                    status_code=503,
                )
            except Exception as e:
                logger.error("Admin data freshness error: %s", e, exc_info=True)
                self._send_json(
                    {"error": f"Freshness check failed: {e}"},
                    status_code=500,
                )

        # ── Data Refresh Status (GET /api/data/refresh-status) ──
        elif path == "/api/data/refresh-status":
            try:
                from data_refresh import get_refresh_status

                status = get_refresh_status()
                self._send_json(status)
            except ImportError:
                self._send_json(
                    {"error": "data_refresh module not available", "running": False},
                    status_code=503,
                )
            except Exception as e:
                logger.error("Refresh status error: %s", e, exc_info=True)
                self._send_json(
                    {"error": f"Refresh status check failed: {e}"},
                    status_code=500,
                )

        # ── ATS Widget embed code (GET /api/widget/embed) ──
        elif path == "/api/widget/embed":
            try:
                from ats_widget import generate_embed_code

                params = urllib.parse.parse_qs(parsed.query)
                widget_cfg = {
                    "apiEndpoint": (params.get("api_endpoint", [""])[0] or "").strip()
                    or "https://media-plan-generator.onrender.com",
                    "jobTitle": (params.get("job_title", [""])[0] or "").strip(),
                    "location": (params.get("location", [""])[0] or "").strip(),
                    "budget": (params.get("budget", ["5000"])[0] or "5000").strip(),
                    "theme": (params.get("theme", ["light"])[0] or "light").strip(),
                    "position": (
                        params.get("position", ["bottom-right"])[0] or "bottom-right"
                    ).strip(),
                }
                embed_code = generate_embed_code(widget_cfg)
                self._send_json(
                    {
                        "embed_code": embed_code,
                        "config": widget_cfg,
                        "version": "1.0.0",
                        "docs": "Add the embed_code snippet to your ATS page HTML to show AI plan suggestions.",
                    }
                )
            except ImportError:
                self._send_json(
                    {"error": "ATS widget module not available"}, status_code=503
                )
            except (ValueError, TypeError, KeyError) as e:
                logger.error(f"Widget embed error: {e}", exc_info=True)
                self._send_json(
                    {"error": f"Failed to generate embed code: {e}"}, status_code=500
                )

        # ── Attribution Dashboard (GET /attribution) ──
        elif path == "/attribution":
            try:
                if not _attribution_available:
                    self.send_error(503, "Attribution module not available")
                    return
                # Use cached report or generate sample
                from attribution_dashboard import _report_cache, _lock as _attr_lock

                with _attr_lock:
                    cached = _report_cache.get("last_report")
                if cached:
                    html = _attribution_html(cached)
                else:
                    # Generate a sample dashboard with demo data
                    sample_data = {
                        "total_budget": 150000,
                        "industry": "technology",
                        "total_hiring_cost": 525000,
                        "channels": [
                            {"channel": "LinkedIn", "spend": 45000},
                            {"channel": "Indeed", "spend": 30000},
                            {"channel": "Google Ads", "spend": 25000},
                            {"channel": "Programmatic", "spend": 20000},
                            {"channel": "Glassdoor", "spend": 15000},
                            {"channel": "Niche Job Boards", "spend": 15000},
                        ],
                    }
                    report = _attribution_report(sample_data)
                    html = _attribution_html(report)
                html_bytes = html.encode("utf-8")
                self._send_compressed_response(html_bytes, "text/html; charset=utf-8")
            except Exception as e:
                logger.error("Attribution dashboard error: %s", e, exc_info=True)
                self.send_error(500, "Attribution dashboard failed")

        # ── Calendar Sync MCP: GET /api/calendar/events ──
        elif path == "/api/calendar/events":
            if not _calendar_sync_available:
                self._send_json(
                    {"error": "Calendar sync not available", "events": []},
                    status_code=503,
                )
                return
            try:
                qs = urllib.parse.parse_qs(urllib.parse.urlparse(self.path).query)
                days = int(qs.get("days", ["30"])[0])
                days = max(1, min(days, 365))
                events = _cal_get_events(days=days)
                self._send_json({"events": events, "count": len(events), "days": days})
            except (ValueError, TypeError) as exc:
                logger.error("Calendar events error: %s", exc, exc_info=True)
                self._send_json(
                    {"error": f"Invalid parameter: {exc}", "events": []},
                    status_code=400,
                )
            except Exception as exc:
                logger.error("Calendar events error: %s", exc, exc_info=True)
                if _slack_alerts_available:
                    _slack_error_alert(exc, context="GET /api/calendar/events")
                self._send_json(
                    {"error": "Calendar service error", "events": []},
                    status_code=500,
                )

        # ── Chroma RAG MCP: GET /api/nova/rag/search?q=... ──
        elif path == "/api/nova/rag/search":
            if not _chroma_rag_available:
                self._send_json(
                    {"error": "ChromaDB RAG not available", "results": []},
                    status_code=503,
                )
                return
            try:
                qs = urllib.parse.parse_qs(urllib.parse.urlparse(self.path).query)
                query = qs.get("q", [""])[0]
                top_k = int(qs.get("top_k", ["5"])[0])
                top_k = max(1, min(top_k, 20))
                if not query:
                    self._send_json(
                        {"error": "Missing 'q' parameter", "results": []},
                        status_code=400,
                    )
                    return
                results = _chroma_search(query, top_k=top_k)
                self._send_json(
                    {"results": results, "count": len(results), "query": query}
                )
            except (ValueError, TypeError) as exc:
                logger.error("Chroma RAG search error: %s", exc, exc_info=True)
                self._send_json(
                    {"error": f"Invalid parameter: {exc}", "results": []},
                    status_code=400,
                )
            except Exception as exc:
                logger.error("Chroma RAG search error: %s", exc, exc_info=True)
                if _slack_alerts_available:
                    _slack_error_alert(exc, context="GET /api/nova/rag/search")
                self._send_json(
                    {"error": "RAG search service error", "results": []},
                    status_code=500,
                )

        # ── Chroma RAG MCP: GET /api/nova/rag/stats ──
        elif path == "/api/nova/rag/stats":
            if not _chroma_rag_available:
                self._send_json({"available": False, "error": "ChromaDB not available"})
                return
            try:
                stats = _chroma_stats()
                self._send_json(stats)
            except Exception as exc:
                logger.error("Chroma RAG stats error: %s", exc, exc_info=True)
                self._send_json({"error": _internal_error_msg()}, status_code=500)

        # ── Nova Intelligent Cache: GET /api/nova/cache/stats ──
        elif path == "/api/nova/cache/stats":
            if not _nova_cache_available:
                self._send_json(
                    {
                        "available": False,
                        "error": "Nova intelligent cache not available",
                    }
                )
                return
            try:
                stats = _nova_cache_stats()
                self._send_json(stats)
            except Exception as exc:
                logger.error("Nova cache stats error: %s", exc, exc_info=True)
                self._send_json({"error": _internal_error_msg()}, status_code=500)

        # ── Slack Alerts MCP: GET /api/slack/alerts/status ──
        elif path == "/api/slack/alerts/status":
            if not _slack_alerts_available:
                self._send_json(
                    {"available": False, "error": "Slack alerts not available"}
                )
                return
            try:
                status = _slack_alerts_status()
                self._send_json(status)
            except Exception as exc:
                logger.error("Slack alerts status error: %s", exc, exc_info=True)
                self._send_json({"error": _internal_error_msg()}, status_code=500)

        # NOTE: /api/pricing/live now handled by routes/pricing.py
        # NOTE: /api/campaign/list, /auto-qc, /eval-framework now handled by
        # routes/campaign.py and routes/pages.py
        else:
            self.send_error(404)

    def do_POST(self):
        # ── Request ID generation (Feature 6) ──
        try:
            from monitoring import generate_request_id as _gen_rid

            self._request_id = _gen_rid()
        except Exception:
            self._request_id = uuid.uuid4().hex[:12]
        # ── Set request context for structured logging ──
        try:
            from monitoring import set_request_id as _set_rid

            _set_rid(self._request_id)
        except ImportError:
            pass
        _req_start = time.time()
        if _metrics:
            _metrics.enter_request()

        # ── Sentry Performance: wrap entire POST in a transaction (SDK 2.x) ──
        parsed_for_sentry = urlparse(self.path)
        _sentry_txn_cm, self._sentry_span = _sentry_start_txn(
            "http.server", f"POST {parsed_for_sentry.path}"
        )

        self._response_status: int = 200
        try:
            with _sentry_txn_cm:
                self._handle_POST()
        except Exception as _exc:
            logger.exception("Unhandled exception in do_POST: %s", _exc)
            self._response_status = 500
            # Slack alert for unhandled POST errors
            if _slack_alerts_available:
                try:
                    _slack_error_alert(
                        _exc,
                        context="Unhandled exception in do_POST",
                        endpoint=urlparse(self.path).path,
                    )
                except Exception:
                    pass
            try:
                self._send_error(
                    "Internal server error",
                    "INTERNAL_ERROR",
                    500,
                )
            except Exception:
                pass
        except BaseException as _bexc:
            # Catch SystemExit / GreenletExit / KeyboardInterrupt -- gevent worker
            # timeout kills greenlets with SystemExit. Log and return a 503 so the
            # worker process stays alive instead of crashing (502).
            logger.error(
                "BaseException in do_POST (worker timeout?): %s", _bexc, exc_info=True
            )
            self._response_status = 503
            try:
                self.send_error(503, "Service Unavailable")
            except Exception:
                pass
        finally:
            _latency = (time.time() - _req_start) * 1000
            if _metrics:
                _metrics.exit_request()
                parsed = urlparse(self.path)
                _metrics.record_request(
                    parsed.path, "POST", self._response_status, _latency
                )
            # ── Slow endpoint profiling ──
            _record_slow_endpoint(urlparse(self.path).path, "POST", _latency)
            # ── SLO tracker + anomaly detector (S17) ──
            try:
                from slo_config import get_slo_tracker

                get_slo_tracker().record_from_endpoint(
                    urlparse(self.path).path, _latency, self._response_status
                )
            except Exception:
                pass
            try:
                from anomaly_detector import record_metric, METRIC_REQUEST_LATENCY

                record_metric(METRIC_REQUEST_LATENCY, _latency)
            except Exception:
                pass
            # ── PostHog: Track download events for any /download/ or /export endpoint ──
            if _posthog_available and self._response_status == 200:
                _ph_post_path = urlparse(self.path).path
                if "/download/" in _ph_post_path or "/export/" in _ph_post_path:
                    try:
                        # Determine format from path
                        _dl_format = "unknown"
                        if "excel" in _ph_post_path or "xlsx" in _ph_post_path:
                            _dl_format = "excel"
                        elif "ppt" in _ph_post_path:
                            _dl_format = "ppt"
                        elif "csv" in _ph_post_path:
                            _dl_format = "csv"
                        elif "html" in _ph_post_path:
                            _dl_format = "html"
                        # Determine product from path
                        _dl_product = (
                            _ph_post_path.split("/api/")[1].split("/")[0]
                            if "/api/" in _ph_post_path
                            else "unknown"
                        )
                        self._ph_track(
                            "report_downloaded",
                            {
                                "format": _dl_format,
                                "product": _dl_product,
                                "path": _ph_post_path,
                            },
                        )
                    except Exception:
                        pass  # Never let analytics break the request
            # ── Clear request context for thread reuse ──
            try:
                from monitoring import clear_request_context as _clear_ctx

                _clear_ctx()
            except ImportError:
                pass

    def _handle_POST(self):
        parsed = urlparse(self.path)
        path = parsed.path

        # ── Idempotency: return cached response for duplicate mutation requests ──
        if self._check_idempotency(path):
            return

        # ── Global body size limits ──
        try:
            content_length = int(self.headers.get("Content-Length") or 0)
        except (ValueError, TypeError):
            content_length = 0

        MAX_BODY_SIZE = 10485760  # 10 MB absolute max (file uploads)
        MAX_API_BODY_SIZE = 1048576  # 1 MB for non-file-upload routes

        # File-upload routes that may need the larger 10 MB limit
        _file_upload_routes = (
            "/api/generate",
            "/api/social-plan/download",
            "/api/audit/analyze",
            "/api/tracker/analyze",
            "/api/hire-signal/analyze",
            "/api/elevenlabs/stt",
        )

        if content_length > MAX_BODY_SIZE:
            self._send_error(
                "Request body too large. Maximum size is 10 MB.",
                "PAYLOAD_TOO_LARGE",
                413,
            )
            return

        if content_length > MAX_API_BODY_SIZE and not any(
            path.startswith(r) for r in _file_upload_routes
        ):
            self._send_error(
                "Request body too large. Maximum size is 1 MB for this endpoint.",
                "PAYLOAD_TOO_LARGE",
                413,
            )
            return

        # ── Bug #16 fix: Validate Content-Type for JSON API endpoints ──
        # Routes that accept non-JSON bodies (file uploads, audio, etc.)
        _non_json_post_routes = (
            "/api/elevenlabs/stt",
            "/api/elevenlabs/sfx",
            "/api/elevenlabs/voiceover",
            "/api/elevenlabs/audio-summary",
            "/api/sentry/webhook",
            "/api/slack/events",
        )
        if (
            path.startswith("/api/")
            and path not in _non_json_post_routes
            and content_length > 0
        ):
            _ct = (self.headers.get("Content-Type") or "").lower()
            if _ct and "application/json" not in _ct and "multipart/" not in _ct:
                self.send_response(415)
                self.send_header("Content-Type", "application/json")
                self.end_headers()
                self.wfile.write(
                    json.dumps(
                        {
                            "error": "Unsupported Media Type. Expected Content-Type: application/json"
                        }
                    ).encode()
                )
                return

        # ── API Versioning: strip /v1 prefix (Feature 4) ──
        if path.startswith("/v1/"):
            path = path[3:]

        # ── API Key Authentication (Phase 6) ──
        if _auth_module_loaded:
            try:
                from auth import authenticate

                auth_header = self.headers.get("Authorization") or ""
                query_params = urllib.parse.parse_qs(
                    urllib.parse.urlparse(self.path).query
                )
                query_key = query_params.get("api_key", [None])[0]
                auth_result = authenticate(path, auth_header, query_key)
                if not auth_result["authenticated"]:
                    self._send_json(
                        {
                            "error": "Authentication required",
                            "reason": auth_result["reason"],
                        },
                        status_code=401,
                    )
                    return
            except Exception:
                logger.error("Auth check failed unexpectedly", exc_info=True)
                self._send_json({"error": "Internal auth error"}, status_code=500)
                return
        elif path.startswith("/api/admin"):
            # Auth module not loaded -- block admin endpoints
            self._send_json(
                {"error": "Auth module unavailable -- admin access denied"},
                status_code=503,
            )
            return

        # ── Audit logging for API access ──
        try:
            from audit_logger import log_event

            if path.startswith("/api/"):
                log_event(
                    "api.access", resource=path, ip_address=self.client_address[0]
                )
        except ImportError:
            pass

        # ── Centralized route-aware rate limiting ──
        # Skip rate limiting for admin routes when valid auth is provided.
        _is_admin = self._check_admin_auth()
        # LLM-heavy analysis endpoints that need stricter limits
        _LLM_HEAVY_PATHS = (
            "/api/competitive/analyze",
            "/api/competitive/scrape",
            "/api/social-plan/generate",
            "/api/talent-heatmap/analyze",
            "/api/audit/analyze",
            "/api/hire-signal/analyze",
            "/api/simulator/simulate",
            "/api/simulator/optimize",
            "/api/compliance/analyze",
            "/api/complianceguard/audit",
            "/api/creative/generate",
            "/api/ab-test/generate",
            "/api/roi/calculate",
        )
        if not _is_admin and path.startswith("/api/"):
            client_ip = self._get_client_ip()
            if path == "/api/generate":
                _rl_allowed = _rl_generate.is_allowed(
                    client_ip, max_requests=10, window_seconds=60
                )
            elif path in ("/api/chat", "/api/chat/stream", "/api/nova/chat"):
                _rl_allowed = _rl_chat.is_allowed(
                    client_ip, max_requests=15, window_seconds=60
                )
            elif path in _LLM_HEAVY_PATHS:
                _rl_allowed = _rl_llm_heavy.is_allowed(
                    client_ip, max_requests=10, window_seconds=60
                )
            elif path.startswith("/api/portal/"):
                _rl_allowed = _rl_portal.is_allowed(
                    client_ip, max_requests=20, window_seconds=60
                )
            else:
                _rl_allowed = _rl_general.is_allowed(
                    client_ip, max_requests=30, window_seconds=60
                )
            if not _rl_allowed:
                self._send_error(
                    "Too many requests. Please try again later.",
                    "RATE_LIMITED",
                    429,
                )
                return

        # ── CSRF token validation (double-submit cookie pattern) ──
        # Skip for API-key-authenticated requests (Bearer token) or admin auth.
        _has_api_key = bool(
            (self.headers.get("Authorization") or "").startswith("Bearer ")
        )
        # Paths exempt from CSRF: external webhooks that use their own signature
        _CSRF_EXEMPT_PATHS = (
            "/api/sentry/webhook",
            "/api/slack/events",
            # S48: Chat endpoints exempt from CSRF -- protected by @joveo.com auth instead.
            # S46 added CSRF to /api/chat but it causes persistent 401s due to:
            # 1. Token expiry (4h TTL) mid-session
            # 2. Race condition on page load (async fetch not awaited)
            # 3. Embedded widget lacks CSRF context
            # The @joveo.com cookie check + rate limiting is sufficient protection.
            "/api/chat",
            "/api/chat/stream",
            "/api/chat/stop",
            "/api/chat/title",
            "/api/chat/feedback",
            "/api/chat/share",
            "/api/generate",  # S48: Protected by same-origin + @joveo.com auth instead
            "/api/health",  # Read-only health checks
            "/api/health/ping",  # Read-only ping
            "/api/csrf-token",  # Must be exempt to bootstrap
            "/api/config",  # Public config, read-only
            "/api/auth/me",  # Uses Bearer token, not CSRF
            "/api/auth/session",  # Uses Bearer token, not CSRF
        )
        if (
            not _has_api_key
            and not _is_admin
            and path.startswith("/api/")
            and path not in _CSRF_EXEMPT_PATHS
        ):
            header_token = self.headers.get("X-CSRF-Token") or ""
            cookie_header = self.headers.get("Cookie") or ""
            cookie_token = _parse_cookie_value(cookie_header, "csrf_token")
            if not _validate_csrf_double_submit(cookie_token, header_token):
                self._send_error(
                    "Invalid or missing CSRF token. Refresh the page and try again.",
                    "AUTH_REQUIRED",
                    403,
                )
                return

        # ── Auth session: mint signed server-side cookie after Supabase login ──
        # S60: Previously this endpoint was fire-and-forget (just logged the
        # event) and the *client JS* wrote the `nova_user_email` cookie
        # unsigned. Under STRICT_AUTH=1 the backend rejects unsigned cookies
        # (see Path 3 in _check_joveo_auth), so no one could stay logged in.
        #
        # Now the server:
        #   1. VERIFIES the Supabase JWT (HMAC over SUPABASE_JWT_SECRET, same
        #      logic as _check_joveo_auth Path 2).
        #   2. Confirms the request body's `email` matches the JWT claim AND
        #      ends with @joveo.com AND the token isn't expired.
        #   3. Mints a signed cookie:
        #          email|hex(HMAC-SHA256(SESSION_SIGNING_SECRET, email))
        #      and returns it via Set-Cookie.
        #
        # Cookie flags chosen deliberately:
        #   HttpOnly       -> blocks JS from reading the cookie, so an XSS
        #                     payload cannot exfiltrate it.
        #   SameSite=Strict-> browser never sends the cookie on cross-site
        #                     requests, which kills CSRF that relies on the
        #                     session cookie traveling on attacker-initiated
        #                     requests.
        #   Secure         -> cookie only travels over HTTPS in production.
        #   Max-Age=86400  -> 24-hour session; the client re-calls this
        #                     endpoint on each Supabase refresh.
        #
        # Backwards compat: if SUPABASE_JWT_SECRET *or* SESSION_SIGNING_SECRET
        # is unset, we log a warning and fall back to the S59 fire-and-forget
        # behaviour (no cookie minted). S59 graceful-migration code in
        # _check_joveo_auth Path 3 already accepts the legacy unsigned cookie
        # (still written by the client) when STRICT_AUTH != 1, so existing
        # users keep working until the operator flips STRICT_AUTH=1.
        if path == "/api/auth/session":
            try:
                content_length = int(self.headers.get("Content-Length") or 0)
                if content_length > 102400:  # 100KB limit
                    self._send_error("Request too large", "PAYLOAD_TOO_LARGE", 413)
                    return
                body = self.rfile.read(content_length) if content_length > 0 else b"{}"
                auth_data = json.loads(body.decode("utf-8"))
                posted_email = (auth_data.get("email") or "").lower().strip()
                access_token = (auth_data.get("access_token") or "").strip()

                jwt_secret = (os.environ.get("SUPABASE_JWT_SECRET") or "").strip()
                session_secret = (
                    os.environ.get("SESSION_SIGNING_SECRET") or ""
                ).strip()

                # ── Fallback path: env vars unset → keep S59 behaviour ──
                if not jwt_secret or not session_secret:
                    logger.warning(
                        "Auth session: cookie minting disabled (SUPABASE_JWT_SECRET "
                        "or SESSION_SIGNING_SECRET unset); falling back to "
                        "fire-and-forget. user=%s email=%s",
                        auth_data.get("user_id", ""),
                        posted_email,
                    )
                    self._send_json({"status": "ok"})
                    return

                # ── JWT verification (same logic as _check_joveo_auth Path 2) ──
                import base64 as _b64

                if not access_token or access_token.count(".") != 2:
                    self._send_json(
                        {"error": "Missing or malformed JWT"}, status_code=401
                    )
                    return
                header_b64, payload_b64, sig_b64 = access_token.split(".")
                signing_input = f"{header_b64}.{payload_b64}".encode("ascii")
                expected_sig = hmac.new(
                    jwt_secret.encode("utf-8"),
                    signing_input,
                    hashlib.sha256,
                ).digest()
                try:
                    actual_sig = _b64.urlsafe_b64decode(
                        sig_b64 + "=" * (4 - len(sig_b64) % 4)
                    )
                except Exception:
                    self._send_json(
                        {"error": "Invalid JWT signature encoding"},
                        status_code=401,
                    )
                    return
                if not hmac.compare_digest(expected_sig, actual_sig):
                    self._send_json({"error": "Invalid JWT signature"}, status_code=401)
                    return

                # Decode payload and verify claims
                try:
                    payload_padded = payload_b64 + "=" * (4 - len(payload_b64) % 4)
                    decoded = json.loads(_b64.urlsafe_b64decode(payload_padded))
                except Exception:
                    self._send_json({"error": "Invalid JWT payload"}, status_code=401)
                    return
                _exp = decoded.get("exp")
                if isinstance(_exp, (int, float)) and _exp < time.time():
                    self._send_json({"error": "JWT expired"}, status_code=401)
                    return
                jwt_email = (decoded.get("email") or "").lower().strip()
                if not jwt_email.endswith("@joveo.com"):
                    self._send_json({"error": "Email not authorized"}, status_code=401)
                    return
                # Body email must match JWT claim — blocks a caller from
                # presenting Alice's token while claiming to be Bob.
                if posted_email and posted_email != jwt_email:
                    self._send_json({"error": "Email / JWT mismatch"}, status_code=401)
                    return

                # ── Mint signed cookie ──
                signing_email = jwt_email
                sig_hex = hmac.new(
                    session_secret.encode("utf-8"),
                    signing_email.encode("utf-8"),
                    hashlib.sha256,
                ).hexdigest()
                signed_value = f"{signing_email}|{sig_hex}"
                is_https = (
                    self.headers.get("X-Forwarded-Proto") or ""
                ).lower() == "https"
                cookie_parts = [
                    f"nova_user_email={signed_value}",
                    "Path=/",
                    "HttpOnly",  # XSS cannot read the cookie
                    "SameSite=Strict",  # CSRF cannot send it cross-site
                    "Max-Age=86400",
                ]
                if is_https:
                    cookie_parts.append("Secure")  # HTTPS-only in production
                cookie_val = "; ".join(cookie_parts)

                # Manual response so we can attach Set-Cookie before body
                resp_body = json.dumps({"status": "ok", "expires_in": 86400}).encode(
                    "utf-8"
                )
                self.send_response(200)
                self.send_header("Content-Type", "application/json")
                self.send_header("Set-Cookie", cookie_val)
                cors_origin = self._get_cors_origin()
                if cors_origin:
                    self.send_header("Access-Control-Allow-Origin", cors_origin)
                    self.send_header("Access-Control-Allow-Credentials", "true")
                self.send_header("Content-Length", str(len(resp_body)))
                self.end_headers()
                self.wfile.write(resp_body)
                logger.info(
                    "Auth session: minted signed cookie for user=%s email=%s",
                    auth_data.get("user_id", ""),
                    signing_email,
                )
            except Exception as e:
                logger.warning("Auth session error: %s", e)
                # Fail closed on unexpected errors (don't leak a cookie
                # after a partial failure).
                try:
                    self._send_json({"error": "Auth session failed"}, status_code=500)
                except Exception:
                    pass
            return

        # ── Extracted POST route modules ──
        from routes.competitive import handle_competitive_post_routes
        from routes.export import handle_export_post_routes
        from routes.copilot import handle_copilot_post_routes
        from routes.campaign import handle_campaign_post_routes
        from routes.pricing import handle_pricing_post_routes
        from routes.tts import handle_tts_post_routes
        from routes.diagram import handle_diagram_post_routes
        from routes.canvas import handle_canvas_post_routes
        from routes.slotops import handle_slotops_post_routes

        if handle_slotops_post_routes(self, path, parsed):
            return
        if handle_tts_post_routes(self, path, parsed):
            return
        if handle_diagram_post_routes(self, path, parsed):
            return
        if handle_competitive_post_routes(self, path, parsed):
            return
        if handle_export_post_routes(self, path, parsed):
            return
        if handle_copilot_post_routes(self, path, parsed):
            return
        if handle_campaign_post_routes(self, path, parsed):
            return
        if handle_pricing_post_routes(self, path, parsed):
            return
        if handle_canvas_post_routes(self, path, parsed):
            return

        # ── Event-sourced plan state machine (plan_events module): undo ──
        if path.startswith("/api/plans/") and path.endswith("/undo"):
            # POST /api/plans/<plan_id>/undo
            parts = path.split("/")
            # Expected: ['', 'api', 'plans', '<plan_id>', 'undo']
            _pe_plan_id = parts[3] if len(parts) >= 5 else ""
            if not _pe_plan_id:
                self._send_json({"error": "Missing plan_id"}, status_code=400)
                return
            try:
                from plan_events import get_event_store as _get_pe_store

                _pe_store = _get_pe_store()
                # Read optional user_id from body
                content_len = int(self.headers.get("Content-Length") or 0)
                if content_len > 102400:  # 100KB limit
                    self._send_error("Request too large", "PAYLOAD_TOO_LARGE", 413)
                    return
                raw = self.rfile.read(content_len) if content_len else b"{}"
                body = json.loads(raw) if raw.strip() else {}
                pe_user = body.get("user_id") or "system"
                compensating = _pe_store.undo(_pe_plan_id, user_id=pe_user)
                self._send_json(
                    {
                        "plan_id": _pe_plan_id,
                        "compensating_event": compensating.to_dict(),
                        "current_snapshot": _pe_store.get_snapshot(_pe_plan_id),
                    }
                )
            except ImportError:
                self._send_json(
                    {"error": "plan_events module not available"}, status_code=503
                )
            except ValueError as ve:
                self._send_json({"error": str(ve)}, status_code=400)
            except Exception as e:
                logger.error(f"Plan events undo error: {e}", exc_info=True)
                self._send_json({"error": _internal_error_msg()}, status_code=500)
            return

        # ── Event-sourced plan state (legacy event_store): undo / redo ──
        if path.startswith("/api/plan/events/") and path.endswith("/undo"):
            # POST /api/plan/events/<plan_id>/undo
            _undo_cl = int(self.headers.get("Content-Length") or 0)
            if _undo_cl > 102400:  # 100KB limit
                self._send_error("Request too large", "PAYLOAD_TOO_LARGE", 413)
                return
            parts = path.split("/")
            # Expected: ['', 'api', 'plan', 'events', '<plan_id>', 'undo']
            _es_plan_id = parts[4] if len(parts) >= 6 else ""
            if not _es_plan_id:
                self._send_json({"error": "Missing plan_id"}, status_code=400)
                return
            try:
                from event_store import get_event_store

                _es = get_event_store()
                result = _es.undo(_es_plan_id)
                status = 400 if "error" in result else 200
                self._send_json(result, status_code=status)
            except ImportError:
                self._send_json(
                    {"error": "event_store module not available"}, status_code=503
                )
            except Exception as e:
                logger.error(f"Event store undo error: {e}", exc_info=True)
                self._send_json({"error": _internal_error_msg()}, status_code=500)
            return
        if path.startswith("/api/plan/events/") and path.endswith("/redo"):
            # POST /api/plan/events/<plan_id>/redo
            _redo_cl = int(self.headers.get("Content-Length") or 0)
            if _redo_cl > 102400:  # 100KB limit
                self._send_error("Request too large", "PAYLOAD_TOO_LARGE", 413)
                return
            parts = path.split("/")
            _es_plan_id = parts[4] if len(parts) >= 6 else ""
            if not _es_plan_id:
                self._send_json({"error": "Missing plan_id"}, status_code=400)
                return
            try:
                from event_store import get_event_store

                _es = get_event_store()
                result = _es.redo(_es_plan_id)
                status = 400 if "error" in result else 200
                self._send_json(result, status_code=status)
            except ImportError:
                self._send_json(
                    {"error": "event_store module not available"}, status_code=503
                )
            except Exception as e:
                logger.error(f"Event store redo error: {e}", exc_info=True)
                self._send_json({"error": _internal_error_msg()}, status_code=500)
            return
        # ── Multi-Agent Negotiation (Linear JOV-23) ──
        if path == "/api/plan/negotiate":
            # ── S46: Server-side @joveo.com domain enforcement ──
            if not self._check_joveo_auth() and not self._check_admin_auth():
                self._send_error(
                    "Authentication required. Please sign in with your @joveo.com account.",
                    "AUTH_REQUIRED",
                    401,
                )
                return
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
                raw = self.rfile.read(content_len) if content_len else b"{}"
                body = json.loads(raw) if raw.strip() else {}
                from multi_agent import handle_negotiate_request

                result, status = handle_negotiate_request(body)
                self._send_json(result, status_code=status)
            except json.JSONDecodeError:
                self._send_json(
                    {"error": "Invalid JSON in request body"}, status_code=400
                )
            except ImportError:
                self._send_json(
                    {"error": "multi_agent module not available"}, status_code=503
                )
            except Exception as e:
                logger.error(f"Plan negotiate error: {e}", exc_info=True)
                self._send_json({"error": _internal_error_msg()}, status_code=500)
            return

        # ── Plan Templates Marketplace: Fork ──
        if path == "/api/templates/fork":
            try:
                from plan_templates import fork_template as _fork_template

                content_len = int(self.headers.get("Content-Length") or 0)
                raw = self.rfile.read(content_len) if content_len else b"{}"
                body = json.loads(raw)
                template_id = body.get("template_id") or ""
                customizations = body.get("customizations") or {}
                if not template_id:
                    self._send_json(
                        {"error": "Missing required field: template_id"},
                        status_code=400,
                    )
                    return
                forked = _fork_template(template_id, customizations)
                if forked is None:
                    self._send_json(
                        {"error": f"Template '{template_id}' not found"},
                        status_code=404,
                    )
                    return
                self._send_json({"forked_template": forked}, status_code=201)
            except ImportError:
                self._send_json(
                    {"error": "Plan templates module not available"}, status_code=503
                )
            except json.JSONDecodeError as e:
                self._send_json({"error": f"Invalid JSON body: {e}"}, status_code=400)
            except (ValueError, TypeError, KeyError) as e:
                logger.error(f"Template fork error: {e}", exc_info=True)
                self._send_json(
                    {"error": f"Failed to fork template: {e}"}, status_code=500
                )
            return

        # ── S47: Save Plan to Supabase (Cindy request) ──
        if path == "/api/saved-plans":
            try:
                if not _check_joveo_auth(self):
                    self._send_error("Authentication required", 401)
                    return
                content_len = int(self.headers.get("Content-Length") or 0)
                body = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body)
                plan_name = (
                    data.get("name") or data.get("plan_name") or "Untitled Plan"
                )[:200]
                plan_data = data.get("data") or data.get("plan_data") or {}
                industry = (data.get("industry") or "")[:100]
                location = (data.get("location") or "")[:200]
                budget = data.get("budget") or 0
                email = _parse_cookie_value(
                    self.headers.get("Cookie") or "", "nova_user_email"
                )
                try:
                    from supabase_client import get_client

                    sb = get_client()
                except Exception:
                    sb = None
                if not sb:
                    self._send_json({"error": "Storage unavailable"}, status_code=503)
                    return
                result = (
                    sb.table("nova_saved_plans")
                    .insert(
                        {
                            "user_email": email or "unknown",
                            "plan_name": plan_name,
                            "plan_data": plan_data,
                            "industry": industry,
                            "location": location,
                            "budget": float(budget),
                        }
                    )
                    .execute()
                )
                saved = (result.data or [{}])[0]
                self._send_json(
                    {"id": saved.get("id"), "plan_name": plan_name, "saved": True}
                )
            except Exception as exc:
                logger.error("Save plan error: %s", exc, exc_info=True)
                self._send_json(
                    {"error": f"Failed to save plan: {exc}"}, status_code=500
                )
            return

        # ── S46: Campaign Optimization API ──
        if path == "/api/optimize":
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
                body = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body)
                from campaign_optimizer import optimize_campaign

                result = optimize_campaign(
                    role=data.get("role") or data.get("job_title") or "",
                    location=data.get("location") or "",
                    industry=data.get("industry") or "",
                    budget=float(data.get("budget") or 0),
                    duration_months=int(data.get("duration_months") or 1),
                    goals=data.get("goals") or [],
                    constraints=data.get("constraints") or {},
                )
                self._send_json(result)
            except json.JSONDecodeError:
                self._send_json(
                    {"error": "Invalid JSON in request body"}, status_code=400
                )
            except Exception as exc:
                logger.error("Campaign optimizer error: %s", exc, exc_info=True)
                self._send_json({"error": "Optimization failed"}, status_code=500)
            return

        if path == "/api/generate":
            # ── S48: Same-origin + @joveo.com auth for plan generation ──
            _gen_auth_ok = (
                self._check_joveo_auth()
                or self._check_admin_auth()
                or "media-plan-generator.onrender.com"
                in (self.headers.get("Origin") or self.headers.get("Referer") or "")
                or "nova.joveo.com"
                in (self.headers.get("Origin") or self.headers.get("Referer") or "")
                or "localhost"
                in (self.headers.get("Origin") or self.headers.get("Referer") or "")
            )
            if not _gen_auth_ok:
                self._send_error(
                    "Authentication required. Please sign in with your @joveo.com account.",
                    "AUTH_REQUIRED",
                    401,
                )
                return
            # ── START TIMEOUT TIMER IMMEDIATELY ──
            # Must cover entire pipeline: validation + enrichment (60s+) + excel generation
            _gen_timeout_flag: threading.Event = threading.Event()
            _gen_timer: threading.Timer = threading.Timer(
                _GENERATE_TIMEOUT_SECONDS, _gen_timeout_flag.set
            )
            _gen_timer.daemon = True
            _gen_timer.start()
            _gen_start: float = time.time()

            if not self._check_rate_limit():
                _gen_timer.cancel()
                self._send_error(
                    "Rate limit exceeded. Please try again in a minute.",
                    "RATE_LIMITED",
                    429,
                )
                return
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
            except (ValueError, TypeError):
                content_len = 0
            if content_len <= 0:
                _gen_timer.cancel()
                self._send_error("Empty request body", "VALIDATION_ERROR", 400)
                return
            if content_len > 10 * 1024 * 1024:  # 10MB limit for JSON API requests
                _gen_timer.cancel()
                self._send_error("Request too large", "VALIDATION_ERROR", 413)
                return
            body = self.rfile.read(content_len)
            try:
                data = json.loads(body)
            except json.JSONDecodeError:
                _gen_timer.cancel()
                self._send_error("Invalid JSON", "VALIDATION_ERROR", 400)
                return

            # ── TIMEOUT CHECK: After JSON parsing ──
            if _gen_timeout_flag.is_set():
                _gen_timer.cancel()
                elapsed: float = time.time() - _gen_start
                logger.error(
                    f"Generation timeout (validation phase) after {elapsed:.1f}s"
                )
                self._send_error(
                    f"Generation timed out after {_GENERATE_TIMEOUT_SECONDS} seconds. Please simplify your request.",
                    "TIMEOUT",
                    504,
                )
                return

            # Validate payload is a JSON object (not array/string/number)
            if not isinstance(data, dict):
                _gen_timer.cancel()
                self._send_error(
                    "Request body must be a JSON object", "VALIDATION_ERROR", 400
                )
                return

            # P3: Max JSON nesting depth (prevent DoS via deeply nested payloads)
            def _check_nesting_depth(
                obj: object, max_depth: int = 5, current: int = 0
            ) -> bool:
                """Return True if nesting depth exceeds max_depth."""
                if current > max_depth:
                    return True
                if isinstance(obj, dict):
                    return any(
                        _check_nesting_depth(v, max_depth, current + 1)
                        for v in obj.values()
                    )
                if isinstance(obj, list):
                    return any(
                        _check_nesting_depth(item, max_depth, current + 1)
                        for item in obj
                    )
                return False

            if _check_nesting_depth(data):
                _gen_timer.cancel()
                self._send_error(
                    "Request JSON nesting too deep (max 5 levels)",
                    "VALIDATION_ERROR",
                    400,
                )
                return

            # ── TIMEOUT CHECK: Before expensive sanitization ──
            if _gen_timeout_flag.is_set():
                _gen_timer.cancel()
                elapsed = time.time() - _gen_start
                logger.error(
                    f"Generation timeout (before sanitization) after {elapsed:.1f}s"
                )
                self._send_error(
                    f"Generation timed out after {_GENERATE_TIMEOUT_SECONDS} seconds. Please simplify your request.",
                    "TIMEOUT",
                    504,
                )
                return

            # Sanitize all string inputs: strip HTML/script tags to prevent stored XSS
            def _sanitize_val(val):
                if val is None:
                    return ""
                if isinstance(val, (int, float, bool)):
                    return str(val)
                if isinstance(val, str):
                    return re.sub(r"<[^>]+>", "", val).strip()
                if isinstance(val, list):
                    return [_sanitize_val(v) for v in val]
                if isinstance(val, dict):
                    return {k: _sanitize_val(v) for k, v in val.items()}
                return val

            for _skey in list(data.keys()):
                data[_skey] = _sanitize_val(data[_skey])

            # Validate required fields
            client_name_input = _safe_str(data.get("client_name")).strip()
            requester_name_input = _safe_str(data.get("requester_name")).strip()
            requester_email_input = _safe_str(data.get("requester_email")).strip()
            roles_input = data.get("target_roles") or data.get("roles") or []
            if isinstance(roles_input, str):
                roles_input = [r.strip() for r in roles_input.split(",") if r.strip()]
            data["target_roles"] = roles_input
            # Bug #2 fix: Auto-wrap string values to lists for array fields
            for _arr_field in (
                "locations",
                "job_categories",
                "competitors",
                "custom_countries",
            ):
                _arr_val = data.get(_arr_field)
                if isinstance(_arr_val, str):
                    data[_arr_field] = [
                        s.strip() for s in _arr_val.split(",") if s.strip()
                    ]
                elif _arr_val is not None and not isinstance(_arr_val, list):
                    data[_arr_field] = [str(_arr_val)]

            # ── S49 FIX (Issue 14): Normalize common location misspellings ──
            # Applied early so all downstream enrichment, PPT, and Excel
            # receive correct city names.
            _LOCATION_CORRECTIONS = {
                "ithica": "Ithaca",
                "harrisburgh": "Harrisburg",
                "murfeesboro": "Murfreesboro",
                "pittsburg": "Pittsburgh",
                "albuquerque": "Albuquerque",
                "cincinatti": "Cincinnati",
                "colombus": "Columbus",
                "detriot": "Detroit",
                "houstan": "Houston",
                "philladelphia": "Philadelphia",
                "phoneix": "Phoenix",
                "sacremento": "Sacramento",
                "seatle": "Seattle",
                "tuscon": "Tucson",
                "milwakee": "Milwaukee",
                "minnapolis": "Minneapolis",
                "indianpolis": "Indianapolis",
                "nashvile": "Nashville",
                "lousville": "Louisville",
                "baltamore": "Baltimore",
                "richmand": "Richmond",
                "charlote": "Charlotte",
                "ralegh": "Raleigh",
                "memhpis": "Memphis",
                "knoxvile": "Knoxville",
                "cleavland": "Cleveland",
                "bufalo": "Buffalo",
                "rochestor": "Rochester",
                "siracuse": "Syracuse",
                "sprinfield": "Springfield",
                "talahassee": "Tallahassee",
                "jackonville": "Jacksonville",
            }

            def _correct_location_spelling(loc_str: str) -> str:
                """Fix common city misspellings while preserving state/country suffix."""
                if not loc_str or not isinstance(loc_str, str):
                    return loc_str
                parts = [p.strip() for p in loc_str.split(",")]
                city_part = parts[0]
                city_lower = city_part.lower().strip()
                corrected = _LOCATION_CORRECTIONS.get(city_lower)
                if corrected:
                    parts[0] = corrected
                    return ", ".join(parts)
                return loc_str

            _locs = data.get("locations") or []
            if isinstance(_locs, list):
                data["locations"] = [
                    _correct_location_spelling(l) if isinstance(l, str) else l
                    for l in _locs
                ]

            # ── S49 FIX (Issue 13): Override "remote" work_environment for
            # trucking/transportation/delivery roles.  OTR trucking is NOT
            # remote work -- drivers are physically on the road.
            _FIELD_WORK_INDUSTRIES = {
                "trucking",
                "transportation",
                "logistics",
                "delivery",
                "freight",
                "shipping",
                "courier",
                "warehousing",
                "supply_chain",
            }
            _FIELD_WORK_ROLES = {
                "driver",
                "trucker",
                "cdl",
                "otr",
                "delivery",
                "courier",
                "forklift",
                "warehouse",
                "dock",
                "loader",
                "freight",
                "hauler",
                "chauffeur",
            }
            _we_raw = _safe_str(data.get("work_environment") or "").strip().lower()
            if _we_raw == "remote":
                _ind_lower = (
                    _safe_str(data.get("industry") or "")
                    .strip()
                    .lower()
                    .replace("-", "_")
                )
                _roles_for_we = data.get("target_roles") or data.get("roles") or []
                if isinstance(_roles_for_we, str):
                    _roles_for_we = [_roles_for_we]
                _roles_text = " ".join(
                    str(r.get("title") if isinstance(r, dict) else r).lower()
                    for r in _roles_for_we
                )
                _industry_is_field = bool(
                    _FIELD_WORK_INDUSTRIES & set(_ind_lower.split("_"))
                )
                _role_is_field = bool(_FIELD_WORK_ROLES & set(_roles_text.split()))
                if _industry_is_field or _role_is_field:
                    data["work_environment"] = "field"
                    data["_work_env_override_note"] = (
                        "Work environment changed from 'remote' to 'field': "
                        "this role involves physical on-site or on-road work."
                    )

            # ── Region selector: validate and normalize ──
            _valid_regions = {"us_only", "global", "emea", "apac", "custom"}
            target_region = (
                _safe_str(data.get("target_region") or "us_only").strip().lower()
            )
            if target_region not in _valid_regions:
                target_region = "us_only"
            data["target_region"] = target_region
            custom_countries = data.get("custom_countries") or []
            if not isinstance(custom_countries, list):
                custom_countries = []
            data["custom_countries"] = [
                c.strip()[:5] for c in custom_countries if isinstance(c, str)
            ][:20]
            if isinstance(roles_input, list) and len(roles_input) > 50:
                self._send_error(
                    "Target roles list exceeds 50 items limit", "VALIDATION_ERROR", 400
                )
                return

            _missing = []
            if not client_name_input:
                _missing.append("Client name")
            if not requester_name_input:
                _missing.append("Requester name")
            if not requester_email_input:
                _missing.append("Requester email")
            if requester_email_input and not _EMAIL_RE.match(requester_email_input):
                _missing.append("Valid email address")
            # Strip any characters that passed regex but are suspicious
            # Must match _EMAIL_RE charset: a-zA-Z0-9._%+-@
            if requester_email_input:
                requester_email_input = re.sub(
                    r"[^a-zA-Z0-9._%+@-]", "", requester_email_input
                )

            # Sanity checks: reject obviously bogus inputs
            if client_name_input and not re.search(r"[a-zA-Z]{2,}", client_name_input):
                _missing.append("Valid client name (must contain letters)")
            if requester_name_input and not re.search(
                r"[a-zA-Z]{2,}", requester_name_input
            ):
                _missing.append("Valid requester name (must contain letters)")

            if not roles_input or (
                isinstance(roles_input, list) and len(roles_input) == 0
            ):
                _missing.append("At least one target role")

            if _missing:
                _gen_timer.cancel()
                self._send_error(
                    f"Required fields missing: {', '.join(_missing)}.",
                    "VALIDATION_ERROR",
                    400,
                )
                return

            # ── TIMEOUT CHECK: Before input length validation ──
            if _gen_timeout_flag.is_set():
                _gen_timer.cancel()
                elapsed = time.time() - _gen_start
                logger.error(
                    f"Generation timeout (after required fields check) after {elapsed:.1f}s"
                )
                self._send_error(
                    f"Generation timed out after {_GENERATE_TIMEOUT_SECONDS} seconds. Please simplify your request.",
                    "TIMEOUT",
                    504,
                )
                return

            # ── API HARDENING: Input length limits (Issue 2) ──
            _use_case = data.get("use_case") or data.get("brief") or ""
            if isinstance(_use_case, str) and len(_use_case) > 10_000:
                self._send_error(
                    "use_case exceeds 10,000 character limit", "VALIDATION_ERROR", 400
                )
                return
            _transcript = data.get("call_transcript") or ""
            if isinstance(_transcript, str) and len(_transcript) > 50_000:
                self._send_error(
                    "call_transcript exceeds 50,000 character limit",
                    "VALIDATION_ERROR",
                    400,
                )
                return
            _competitors = data.get("competitors") or []
            if isinstance(_competitors, list) and len(_competitors) > 20:
                self._send_error(
                    "competitors list exceeds 20 items limit", "VALIDATION_ERROR", 400
                )
                return
            if isinstance(client_name_input, str) and len(client_name_input) > 200:
                self._send_error(
                    "client_name exceeds 200 character limit", "VALIDATION_ERROR", 400
                )
                return

            # ── API HARDENING: Request deduplication (Issue 3) ──
            if _is_duplicate_request(data):
                self._send_error(
                    "Your plan is already being generated. Please wait for the current request to complete, or try again in 60 seconds.",
                    "RATE_LIMITED",
                    429,
                )
                return

            data["client_name"] = client_name_input
            data["requester_name"] = requester_name_input
            data["requester_email"] = requester_email_input

            # Validate critical input fields (non-blocking: store warnings, don't 400)
            _roles_input = data.get("target_roles") or data.get("roles") or []
            _locs_input = data.get("locations") or []
            if isinstance(_locs_input, list) and len(_locs_input) > 100:
                self._send_error(
                    "Locations list exceeds 100 items limit", "VALIDATION_ERROR", 400
                )
                return
            if isinstance(_roles_input, list) and len(_roles_input) > 30:
                self._send_error(
                    "Maximum 30 roles supported per plan. Please split into multiple plans.",
                    "VALIDATION_ERROR",
                    400,
                )
                return

            # P2: Numeric range validation -- budget cap
            _budget_raw_str = _safe_str(
                data.get("budget") or data.get("budget_range") or ""
            )
            _bval = parse_budget(_budget_raw_str)
            if parse_budget.last_was_defaulted and _budget_raw_str.strip():
                self._send_error(
                    f"Could not parse budget '{_budget_raw_str}'. "
                    "Please enter a numeric value like $50,000",
                    "VALIDATION_ERROR",
                    400,
                )
                return
            if _bval <= 0:
                _bval = 100_000.0  # Default to $100K instead of blocking
            if (
                _bval > 1_000_000_000
            ):  # $1B cap (raised from $100M to support large enterprise budgets)
                self._send_error(
                    "Budget exceeds maximum allowed value ($1B). Please enter a value under $1,000,000,000.",
                    "VALIDATION_ERROR",
                    400,
                )
                return

            # ── CRITICAL: Validate budget is explicitly set (Ashlie Issue #1) ──
            _budget_input = str(
                data.get("budget") or "" or data.get("budget_range") or "" or ""
            ).strip()
            if not _budget_input or _budget_input == "":
                self._send_error(
                    "Budget must be specified. Please select a budget range or enter an exact amount.",
                    "VALIDATION_ERROR",
                    400,
                )
                return

            # ── CRITICAL: Validate locations are valid (Ashlie Issue #2) ──
            _invalid_locations = ["hell", "test", "xxx", "fake", "placeholder", "n/a"]
            if isinstance(_locs_input, list):
                for loc in _locs_input:
                    loc_lower = str(loc or "").lower().strip()
                    if loc_lower in _invalid_locations:
                        self._send_error(
                            f'Location "{loc}" is not valid. Please enter a real city, state, or country.',
                            "VALIDATION_ERROR",
                            400,
                        )
                        return
                    # Reject locations that are too short or have no letters
                    if len(loc_lower) < 2 or not any(c.isalpha() for c in loc_lower):
                        self._send_error(
                            f'Location "{loc}" is invalid. Please use a real city, state, or country name.',
                            "VALIDATION_ERROR",
                            400,
                        )
                        return

            # ── CRITICAL: Validate hire_volume is reasonable (Ashlie Issue #3) ──
            _hire_vol = data.get("hire_volume") or ""
            if _hire_vol:
                # Extract number from strings like "100 hires" or "100-500"
                import re as _re_module

                _hire_match = _re_module.search(r"(\d+)", str(_hire_vol))
                if _hire_match:
                    _hire_num = int(_hire_match.group(1))
                    if _hire_num <= 0:
                        self._send_error(
                            "Hire volume must be greater than 0.",
                            "VALIDATION_ERROR",
                            400,
                        )
                        return
                    if _hire_num > 100000:
                        self._send_error(
                            "Hire volume exceeds maximum (100,000).",
                            "VALIDATION_ERROR",
                            400,
                        )
                        return

            # ── Gold Standard: Negative budget rejection ──
            _budget_raw_str = str(
                data.get("budget") or data.get("budget_range") or ""
            ).strip()
            if _budget_raw_str and re.match(
                r"^-", _budget_raw_str.replace("$", "").strip()
            ):
                _gen_timer.cancel()
                self._send_error(
                    "Budget cannot be negative. Please enter a positive value.",
                    "VALIDATION_ERROR",
                    400,
                )
                return

            # ── Budget period normalization (monthly/quarterly/annual → campaign total) ──
            _budget_period = (
                str(data.get("budget_period") or "campaign").strip().lower()
            )
            if _budget_period in ("monthly", "quarterly", "annual"):
                _dur_str = str(data.get("campaign_duration") or "").lower()
                # Extract campaign duration in months from strings like "3 months", "6 months", "1 year"
                _dur_months = 1
                _dur_match = re.search(r"(\d+)", _dur_str)
                if _dur_match:
                    _dur_num = int(_dur_match.group(1))
                    if "year" in _dur_str:
                        _dur_months = _dur_num * 12
                    elif "quarter" in _dur_str:
                        _dur_months = _dur_num * 3
                    else:
                        _dur_months = max(_dur_num, 1)
                # Compute multiplier: how many periods fit in the campaign duration
                _period_months = {"monthly": 1, "quarterly": 3, "annual": 12}[
                    _budget_period
                ]
                _multiplier = max(_dur_months / _period_months, 1.0)
                # Scale the parsed budget value to campaign total
                _bval_period = parse_budget(
                    _safe_str(data.get("budget") or data.get("budget_range") or "")
                )
                if _bval_period > 0:
                    _scaled = _bval_period * _multiplier
                    data["budget"] = f"${_scaled:,.0f}"
                    data["budget_range"] = data["budget"]
                    data["_budget_period_original"] = _budget_period
                    data["_budget_multiplier"] = _multiplier
                    logger.info(
                        f"Budget period normalization: {_budget_period} ${_bval_period:,.0f} x {_multiplier:.1f} = ${_scaled:,.0f}"
                    )

            # ── Gold Standard: Campaign start month validation ──
            _csm_raw = data.get("campaign_start_month")
            if _csm_raw is not None and str(_csm_raw).strip():
                try:
                    _csm_val = int(str(_csm_raw).strip())
                    if _csm_val == 0:
                        # 0 means "use current month" -- auto-set
                        import datetime as _dt_fix

                        data["campaign_start_month"] = _dt_fix.datetime.now().month
                    elif _csm_val < 1 or _csm_val > 12:
                        _gen_timer.cancel()
                        self._send_error(
                            f"Invalid campaign_start_month: {_csm_val}. Must be 1-12.",
                            "VALIDATION_ERROR",
                            400,
                        )
                        return
                except (ValueError, TypeError):
                    # Non-numeric -- default to current month silently
                    import datetime as _dt_fix

                    data["campaign_start_month"] = _dt_fix.datetime.now().month

            # ── Gold Standard: Reject javascript:/file:// URLs in all string fields ──
            _DANGEROUS_URL_RE = re.compile(
                r"(?:javascript|file|data|vbscript)\s*:", re.IGNORECASE
            )
            for _fkey, _fval in data.items():
                if isinstance(_fval, str) and _DANGEROUS_URL_RE.search(_fval):
                    _gen_timer.cancel()
                    self._send_error(
                        f"Field '{_fkey}' contains a disallowed URL scheme (javascript/file/data).",
                        "VALIDATION_ERROR",
                        400,
                    )
                    return
                if isinstance(_fval, list):
                    for _item in _fval:
                        if isinstance(_item, str) and _DANGEROUS_URL_RE.search(_item):
                            _gen_timer.cancel()
                            self._send_error(
                                f"Field '{_fkey}' contains a disallowed URL scheme.",
                                "VALIDATION_ERROR",
                                400,
                            )
                            return

            _validation_warnings = []
            if not _roles_input or (
                isinstance(_roles_input, list) and len(_roles_input) == 0
            ):
                _validation_warnings.append(
                    "No target roles specified -- plan will use generic benchmarks"
                )
            if not _locs_input or (
                isinstance(_locs_input, list) and len(_locs_input) == 0
            ):
                _validation_warnings.append(
                    "No locations specified -- plan will use global averages"
                )
            if not _budget_input:
                _validation_warnings.append(
                    "No budget specified -- plan will use estimated ranges"
                )
            if _validation_warnings:
                data["_input_warnings"] = _validation_warnings
                logger.info("Input validation warnings: %s", _validation_warnings)

            # ── Feature 2b: Async generation mode ──
            if (self.headers.get("X-Async") or "").lower() == "true":
                data["_gen_created"] = (
                    time.time()
                )  # Stamp for Slack notification timing
                job_id = uuid.uuid4().hex[:12]
                # Bug #14 fix: Store session token with job for IDOR prevention
                # C1 fix: Use nova_session (HttpOnly) cookie instead of csrf_token
                _job_csrf = _parse_cookie_value(
                    self.headers.get("Cookie") or "", "nova_session"
                ) or _parse_cookie_value(
                    self.headers.get("Cookie") or "", "csrf_token"
                )  # Fallback for pre-existing sessions
                with _generation_jobs_lock:
                    _generation_jobs[job_id] = {
                        "status": "processing",
                        "progress_pct": 0,
                        "status_message": "Starting generation...",
                        "created": time.time(),
                        "result_bytes": None,
                        "result_content_type": None,
                        "result_filename": None,
                        "error": None,
                        "_session_token": _job_csrf,
                    }
                request_id = getattr(self, "_request_id", None)

                def _async_generate(jid, gen_data, rid):
                    """Run the full sync generation pipeline in a background thread."""
                    try:
                        # S29 v2: Quality-first -- give the pipeline enough time
                        # to use ALL data sources (15+ APIs, 25+ KB files, Supabase,
                        # vector search, LLM synthesis). Reliability > speed.
                        _roles_count = len(
                            gen_data.get("target_roles") or gen_data.get("roles") or []
                        )
                        _locs_count = len(gen_data.get("locations") or [])
                        _is_complex = _roles_count >= 5 or _locs_count >= 5
                        _gen_deadline_secs = (
                            120 if _is_complex else 90
                        )  # 1.5-2 min (was 7-8 min)
                        _gen_deadline = time.time() + _gen_deadline_secs
                        # Tight stage timeouts -- fail fast, serve partial data
                        _enrich_timeout = 20  # Parallel APIs -- 20s is plenty (was 60)
                        _synth_timeout = (
                            45  # KB + vector search + LLM narrative (was 90)
                        )
                        _gs_timeout = 20  # Gold standard quality gates (was 45)
                        _verify_timeout = 15  # LLM verification (was 45)
                        _excel_timeout = (
                            45  # Excel with LLM executive narrative (was 90)
                        )
                        _is_simple = False  # deprecated -- always use full pipeline
                        logger.info(
                            "Async generate job %s: roles=%d, locs=%d, simple=%s, deadline=%ds",
                            jid,
                            _roles_count,
                            _locs_count,
                            _is_simple,
                            _gen_deadline_secs,
                        )

                        # PERF: Pipeline timing instrumentation for profiling
                        _pipeline_t0 = time.time()
                        _step_timings: dict = {}

                        with _generation_jobs_lock:
                            if jid in _generation_jobs:
                                _generation_jobs[jid]["progress_pct"] = 10
                                _generation_jobs[jid][
                                    "status_message"
                                ] = "Enriching market data..."

                        # ── TIMEOUT CHECK: Before enrichment (expensive) ──
                        if time.time() > _gen_deadline:
                            raise TimeoutError(
                                "Plan generation exceeded 2-minute time limit"
                            )

                        # Enrichment -- ThreadPoolExecutor with 25s hard timeout.
                        # CRITICAL: Do NOT use context manager (`with ... as pool:`)
                        # because pool.shutdown(wait=True) blocks until thread completes,
                        # defeating the timeout. Instead, create pool, submit, get result,
                        # and shutdown(wait=False) to avoid blocking on semaphore.
                        enriched = {}
                        _t_enrich = time.time()
                        if enrich_data is not None:
                            _enrich_pool = ThreadPoolExecutor(
                                max_workers=1, thread_name_prefix="async-enrich"
                            )
                            try:
                                _enrich_future = (
                                    _enrich_pool.submit(
                                        enrich_data, gen_data, request_id=rid
                                    )
                                    if rid
                                    else _enrich_pool.submit(enrich_data, gen_data)
                                )
                                enriched = (
                                    _enrich_future.result(timeout=_enrich_timeout) or {}
                                )
                                gen_data["_enriched"] = enriched
                            except TimeoutError:
                                logger.error(
                                    "Async enrichment timed out after %ds for job %s -- continuing with empty data",
                                    _enrich_timeout,
                                    jid,
                                )
                                _enrich_future.cancel()
                                gen_data["_enriched"] = {}
                            except Exception:
                                logger.error(
                                    "Async enrichment failed for job %s",
                                    jid,
                                    exc_info=True,
                                )
                                gen_data["_enriched"] = {}
                            finally:
                                # shutdown(wait=False) does NOT block on running threads
                                _enrich_pool.shutdown(wait=False, cancel_futures=True)
                        else:
                            gen_data["_enriched"] = {}

                        _step_timings["enrichment"] = round(time.time() - _t_enrich, 2)
                        logger.info(
                            "PERF: Enrichment took %.2fs", _step_timings["enrichment"]
                        )

                        # ── Fix 25 + S49 Issue 17: Detect enrichment quality degradation ──
                        # Scale warning text with confidence level so PPT disclaimer
                        # is stronger when data quality is worse.
                        _enr = gen_data.get("_enriched") or {}
                        _enr_summary = (
                            _enr.get("enrichment_summary", {})
                            if isinstance(_enr, dict)
                            else {}
                        )
                        _enr_succeeded = _enr_summary.get("apis_succeeded") or []
                        _enr_confidence = float(
                            _enr_summary.get("confidence_score", 0)
                            if isinstance(_enr_summary, dict)
                            else 0
                        )
                        if not _enr or (
                            isinstance(_enr_succeeded, list) and len(_enr_succeeded) < 3
                        ):
                            # S49: Tiered warning based on confidence score
                            if _enr_confidence < 0.40:
                                gen_data["quality_warning"] = (
                                    "Minimal data available. Plan is heavily estimated. "
                                    "Recommend API enrichment before executing."
                                )
                            elif _enr_confidence < 0.60:
                                gen_data["quality_warning"] = (
                                    "Limited data available. This plan uses estimated benchmarks. "
                                    "Results may vary significantly."
                                )
                            else:
                                gen_data["quality_warning"] = (
                                    "Some market data sources were unavailable. "
                                    "Plan uses benchmark estimates."
                                )
                            logger.warning(
                                "Enrichment quality degraded for job %s: %d APIs succeeded (min 3), confidence=%.2f",
                                jid,
                                (
                                    len(_enr_succeeded)
                                    if isinstance(_enr_succeeded, list)
                                    else 0
                                ),
                                _enr_confidence,
                            )

                        with _generation_jobs_lock:
                            if jid in _generation_jobs:
                                _generation_jobs[jid]["progress_pct"] = 30
                                _generation_jobs[jid][
                                    "status_message"
                                ] = "Synthesizing knowledge base..."

                        # KB + Synthesis -- all wrapped in 45s timeout to prevent
                        # Supabase or LLM calls from hanging indefinitely
                        kb = load_knowledge_base()

                        def _supabase_and_synthesize():
                            """Supabase enrichment + vector search + LLM synthesis.

                            S29 v2: Now queries ALL Supabase tables and runs vector
                            search for quality-first plan generation.
                            """
                            _kb = kb
                            _sb_industry = gen_data.get("industry") or ""
                            _roles_raw = (
                                gen_data.get("target_roles")
                                or gen_data.get("roles")
                                or []
                            )
                            _roles_list = [
                                r.get("title") if isinstance(r, dict) else r
                                for r in (
                                    _roles_raw if isinstance(_roles_raw, list) else []
                                )
                            ]
                            _locs_raw = gen_data.get("locations") or []
                            _locs_list = [
                                (
                                    (l.get("city", "") + ", " + l.get("state", ""))
                                    if isinstance(l, dict)
                                    else l
                                )
                                for l in (
                                    _locs_raw if isinstance(_locs_raw, list) else []
                                )
                            ]
                            _first_role = _roles_list[0] if _roles_list else ""
                            _first_loc = _locs_list[0] if _locs_list else ""

                            if _supabase_data_available:
                                try:
                                    # ── Channel benchmarks ──
                                    _sb_benchmarks = get_channel_benchmarks(
                                        industry=_sb_industry
                                    )
                                    if _sb_benchmarks:
                                        _kb["_supabase_channel_benchmarks"] = (
                                            _sb_benchmarks
                                        )

                                    # ── Industry insights ──
                                    _sb_kb_insights = get_knowledge(
                                        "industry_insights", _sb_industry
                                    )
                                    if _sb_kb_insights:
                                        _kb["_supabase_industry_insights"] = (
                                            _sb_kb_insights
                                        )

                                    # ── Vendor profiles ──
                                    try:
                                        _sb_vendors = get_vendor_profiles()
                                        if _sb_vendors:
                                            _kb["_supabase_vendor_profiles"] = (
                                                _sb_vendors
                                            )
                                            gen_data["_vendor_profiles"] = _sb_vendors
                                    except Exception as _vp_err:
                                        logger.warning(
                                            "Vendor profiles failed (non-fatal): %s",
                                            _vp_err,
                                        )

                                    # ── S29: Supply repository (20 publishers) ──
                                    try:
                                        _sb_supply = get_supply_repository()
                                        if _sb_supply:
                                            _kb["_supabase_supply_repository"] = (
                                                _sb_supply
                                            )
                                            gen_data["_supply_repository"] = _sb_supply
                                            logger.info(
                                                "Supabase supply_repository: %d publishers loaded",
                                                len(_sb_supply),
                                            )
                                    except Exception as _sr_err:
                                        logger.warning(
                                            "Supply repository failed (non-fatal): %s",
                                            _sr_err,
                                        )

                                    # ── S30: Expanded supply repository (2.7MB) ──
                                    if _expanded_supply_repo:
                                        gen_data["_expanded_supply_repo"] = (
                                            _expanded_supply_repo
                                        )
                                        logger.info(
                                            "Expanded supply repo: %d entries injected",
                                            (
                                                len(_expanded_supply_repo)
                                                if isinstance(
                                                    _expanded_supply_repo, (list, dict)
                                                )
                                                else 0
                                            ),
                                        )

                                    # ── S30: Client plan references (RTX + others) ──
                                    if _client_plans_data:
                                        _industry_key = (
                                            gen_data.get("industry") or ""
                                        ).lower()
                                        _matching_plans = {}
                                        for pkey, pdata in _client_plans_data.items():
                                            p_industry = (
                                                pdata.get("industry") or ""
                                            ).lower()
                                            if (
                                                _industry_key
                                                and p_industry
                                                and (
                                                    _industry_key in p_industry
                                                    or p_industry in _industry_key
                                                )
                                            ):
                                                _matching_plans[pkey] = pdata
                                        if _matching_plans:
                                            gen_data["_reference_client_plans"] = (
                                                _matching_plans
                                            )
                                            logger.info(
                                                "Matched %d reference client plans for industry '%s': %s",
                                                len(_matching_plans),
                                                _industry_key,
                                                ", ".join(_matching_plans.keys()),
                                            )
                                        else:
                                            # Include all plans as general reference
                                            gen_data["_reference_client_plans"] = (
                                                _client_plans_data
                                            )
                                            logger.info(
                                                "No industry match -- all %d client plans available as reference",
                                                len(_client_plans_data),
                                            )

                                    # ── S29: Salary data (48 rows) ──
                                    try:
                                        _sb_salaries = get_salary_data(
                                            role=_first_role, location=_first_loc
                                        )
                                        if _sb_salaries:
                                            _kb["_supabase_salary_data"] = _sb_salaries
                                            logger.info(
                                                "Supabase salary_data: %d rows for %s/%s",
                                                len(_sb_salaries),
                                                _first_role,
                                                _first_loc,
                                            )
                                    except Exception as _sal_err:
                                        logger.warning(
                                            "Salary data failed (non-fatal): %s",
                                            _sal_err,
                                        )

                                    # ── S29: Market trends (8 rows) ──
                                    try:
                                        _sb_trends = get_market_trends(limit=20)
                                        if _sb_trends:
                                            _kb["_supabase_market_trends"] = _sb_trends
                                            logger.info(
                                                "Supabase market_trends: %d trends loaded",
                                                len(_sb_trends),
                                            )
                                    except Exception as _mt_err:
                                        logger.warning(
                                            "Market trends failed (non-fatal): %s",
                                            _mt_err,
                                        )

                                    # ── S29: Compliance rules (8 rules) ──
                                    try:
                                        _sb_compliance = get_compliance_rules(
                                            jurisdiction=_first_loc
                                        )
                                        if _sb_compliance:
                                            _kb["_supabase_compliance_rules"] = (
                                                _sb_compliance
                                            )
                                            logger.info(
                                                "Supabase compliance_rules: %d rules for %s",
                                                len(_sb_compliance),
                                                _first_loc,
                                            )
                                    except Exception as _cr_err:
                                        logger.warning(
                                            "Compliance rules failed (non-fatal): %s",
                                            _cr_err,
                                        )

                                except (
                                    urllib.error.URLError,
                                    OSError,
                                    ValueError,
                                ) as sb_err:
                                    logger.error(
                                        "Supabase enrichment failed: %s",
                                        sb_err,
                                        exc_info=True,
                                    )

                            # ── S29: Vector search for relevant KB passages ──
                            if _vector_search_available and _vector_search:
                                try:
                                    _vs_query = f"{_first_role} {_first_loc} {_sb_industry} recruitment media plan"
                                    _vs_results = _vector_search(_vs_query, top_k=10)
                                    if _vs_results:
                                        _kb["_vector_search_results"] = _vs_results
                                        logger.info(
                                            "Vector search: %d relevant passages for '%s'",
                                            len(_vs_results),
                                            _vs_query[:60],
                                        )
                                except Exception as _vs_err:
                                    logger.warning(
                                        "Vector search failed (non-fatal): %s", _vs_err
                                    )

                            # ── S29: Load competitor careers data if available ──
                            try:
                                import pathlib as _pathlib

                                _cc_path = (
                                    _pathlib.Path(__file__).parent
                                    / "data"
                                    / "competitor_careers.json"
                                )
                                if _cc_path.exists():
                                    with open(_cc_path, "r") as _cc_f:
                                        _cc_data = json.loads(_cc_f.read())
                                    if _cc_data:
                                        _kb["_competitor_careers"] = _cc_data
                                        logger.info(
                                            "Loaded competitor_careers.json (%d entries)",
                                            (
                                                len(_cc_data)
                                                if isinstance(_cc_data, list)
                                                else 1
                                            ),
                                        )
                            except Exception as _cc_err:
                                logger.warning(
                                    "competitor_careers.json failed (non-fatal): %s",
                                    _cc_err,
                                )

                            gen_data["_knowledge_base"] = _kb
                            if data_synthesize is not None:
                                return data_synthesize(enriched, _kb, gen_data)
                            return {}

                        _t_synth = time.time()
                        _synth_pool = ThreadPoolExecutor(
                            max_workers=1, thread_name_prefix="async-synth"
                        )
                        try:
                            _synth_future = _synth_pool.submit(_supabase_and_synthesize)
                            gen_data["_synthesized"] = (
                                _synth_future.result(timeout=_synth_timeout) or {}
                            )
                        except TimeoutError:
                            logger.warning(
                                "KB synthesis timed out after %ds for job %s -- continuing with empty synthesis",
                                _synth_timeout,
                                jid,
                            )
                            _synth_future.cancel()
                            gen_data["_synthesized"] = {}
                            gen_data["_knowledge_base"] = kb
                        except Exception:
                            gen_data["_synthesized"] = {}
                            gen_data["_knowledge_base"] = kb
                        finally:
                            _synth_pool.shutdown(wait=False, cancel_futures=True)
                        _step_timings["synthesis"] = round(time.time() - _t_synth, 2)
                        logger.info(
                            "PERF: Synthesis took %.2fs", _step_timings["synthesis"]
                        )

                        # Copy geopolitical context from enriched to synthesized
                        # (PPT/Excel read from _synthesized; enrichment stores it in _enriched)
                        if isinstance(enriched, dict) and isinstance(
                            gen_data.get("_synthesized"), dict
                        ):
                            _geo_from_enriched = enriched.get("geopolitical_context")
                            if _geo_from_enriched and isinstance(
                                _geo_from_enriched, dict
                            ):
                                gen_data["_synthesized"][
                                    "geopolitical_context"
                                ] = _geo_from_enriched

                        if time.time() > _gen_deadline:
                            raise TimeoutError(
                                "Plan generation exceeded 2-minute time limit"
                            )

                        with _generation_jobs_lock:
                            if jid in _generation_jobs:
                                _generation_jobs[jid]["progress_pct"] = 40
                                _generation_jobs[jid][
                                    "status_message"
                                ] = "Analyzing channels..."

                        # Industry classification
                        industry_raw = gen_data.get("industry") or ""
                        company_name = gen_data.get("client_name") or ""
                        roles_list = gen_data.get("target_roles") or gen_data.get(
                            "roles", []
                        )
                        if isinstance(roles_list, str):
                            roles_list = [
                                r.strip() for r in roles_list.split(",") if r.strip()
                            ]
                        industry_profile = classify_industry(
                            industry_raw, company_name, roles_list
                        )
                        gen_data["industry"] = industry_profile.get(
                            "legacy_key", "general_entry_level"
                        )
                        if not gen_data.get("industry_label"):
                            gen_data["industry_label"] = industry_profile["sector"]
                        gen_data["talent_profile"] = industry_profile["talent_profile"]
                        gen_data["bls_sector"] = industry_profile["bls_sector"]
                        gen_data["naics_code"] = industry_profile.get("naics", "00")

                        if time.time() > _gen_deadline:
                            raise TimeoutError(
                                "Plan generation exceeded 2-minute time limit"
                            )

                        with _generation_jobs_lock:
                            if jid in _generation_jobs:
                                _generation_jobs[jid]["progress_pct"] = 60
                                _generation_jobs[jid][
                                    "status_message"
                                ] = "Computing budget allocation..."

                        # Budget Allocation (Phase 4 -- same as sync path)
                        _t_budget = time.time()
                        if calculate_budget_allocation is not None:
                            try:
                                _ind_key_ba = gen_data.get(
                                    "industry", "general_entry_level"
                                )
                                _DEFAULT_ALLOC_BA = {
                                    "programmatic_dsp": 35,
                                    "global_boards": 20,
                                    "niche_boards": 15,
                                    "social_media": 12,
                                    "regional_boards": 8,
                                    "employer_branding": 5,
                                    "apac_regional": 3,
                                    "emea_regional": 2,
                                }
                                # CRITICAL: dict.copy() to avoid mutating shared INDUSTRY_ALLOC_PROFILES
                                if INDUSTRY_ALLOC_PROFILES is not None:
                                    channel_pcts = dict(
                                        INDUSTRY_ALLOC_PROFILES.get(
                                            _ind_key_ba, _DEFAULT_ALLOC_BA
                                        )
                                    )
                                else:
                                    channel_pcts = dict(_DEFAULT_ALLOC_BA)

                                # ── Strip APAC/EMEA channels based on target_region (async path) ──
                                _target_region_async = (
                                    gen_data.get("target_region") or "us_only"
                                )
                                _locs_raw_async = gen_data.get("locations") or []
                                _all_us_async = _target_region_async == "us_only"
                                if not _all_us_async and _target_region_async not in (
                                    "global",
                                    "emea",
                                    "apac",
                                    "custom",
                                ):
                                    # Infer from locations if region not explicitly set
                                    _all_us_async = True
                                    if _locs_raw_async:
                                        for _loc_a in (
                                            _locs_raw_async
                                            if isinstance(_locs_raw_async, list)
                                            else [_locs_raw_async]
                                        ):
                                            _loc_str_a = str(
                                                _loc_a.get("country")
                                                if isinstance(_loc_a, dict)
                                                else _loc_a
                                            ).lower()
                                            if _loc_str_a and _loc_str_a not in (
                                                "us",
                                                "usa",
                                                "united states",
                                                "",
                                            ):
                                                _all_us_async = False
                                                break

                                if isinstance(channel_pcts, dict):
                                    if (
                                        _all_us_async
                                        or _target_region_async == "us_only"
                                    ):
                                        _intl_pct_a = channel_pcts.pop(
                                            "apac_regional", 0
                                        ) + channel_pcts.pop("emea_regional", 0)
                                        if _intl_pct_a > 0:
                                            _top_ch_a = max(
                                                channel_pcts,
                                                key=lambda k: channel_pcts[k],
                                            )
                                            channel_pcts[_top_ch_a] = (
                                                channel_pcts[_top_ch_a] + _intl_pct_a
                                            )
                                            logger.info(
                                                "Async US-only plan (region=%s): redistributed %d%% from APAC/EMEA to %s",
                                                _target_region_async,
                                                _intl_pct_a,
                                                _top_ch_a,
                                            )
                                    elif _target_region_async == "emea":
                                        _apac_pct_a = channel_pcts.pop(
                                            "apac_regional", 0
                                        )
                                        if _apac_pct_a > 0:
                                            channel_pcts["emea_regional"] = (
                                                channel_pcts.get("emea_regional", 0)
                                                + _apac_pct_a
                                            )
                                    elif _target_region_async == "apac":
                                        _emea_pct_a = channel_pcts.pop(
                                            "emea_regional", 0
                                        )
                                        if _emea_pct_a > 0:
                                            channel_pcts["apac_regional"] = (
                                                channel_pcts.get("apac_regional", 0)
                                                + _emea_pct_a
                                            )

                                _bstr_ba = str(
                                    gen_data.get("budget")
                                    or ""
                                    or gen_data.get("budget_range")
                                    or ""
                                    or ""
                                ).strip()
                                if not _bstr_ba:
                                    _bstr_ba = str(
                                        gen_data.get("budget_range") or "" or ""
                                    ).strip()
                                _bval_ba = parse_budget(_bstr_ba)

                                _roles_raw = (
                                    gen_data.get("target_roles")
                                    or gen_data.get("roles")
                                    or []
                                )
                                _roles_for_ba = []
                                for r in (
                                    _roles_raw if isinstance(_roles_raw, list) else []
                                ):
                                    if isinstance(r, str):
                                        _roles_for_ba.append(
                                            {
                                                "title": r,
                                                "count": 1,
                                                "tier": gen_data.get("_role_tiers", {})
                                                .get(r, {})
                                                .get("tier", "Professional"),
                                            }
                                        )
                                    elif isinstance(r, dict):
                                        _roles_for_ba.append(
                                            {
                                                "title": r.get("title") or "",
                                                "count": int(r.get("count", 1)),
                                                "tier": r.get("_tier", "Professional"),
                                            }
                                        )

                                _locs_raw = gen_data.get("locations") or []
                                _locs_for_ba = []
                                for loc in (
                                    _locs_raw if isinstance(_locs_raw, list) else []
                                ):
                                    if isinstance(loc, str):
                                        parts = [p.strip() for p in loc.split(",")]
                                        _locs_for_ba.append(
                                            {
                                                "city": parts[0] if parts else loc,
                                                "state": (
                                                    parts[1] if len(parts) > 1 else ""
                                                ),
                                                "country": (
                                                    parts[2]
                                                    if len(parts) > 2
                                                    else (
                                                        parts[-1]
                                                        if len(parts) > 1
                                                        else "US"
                                                    )
                                                ),
                                            }
                                        )
                                    elif isinstance(loc, dict):
                                        _locs_for_ba.append(
                                            {
                                                "city": loc.get("city") or "",
                                                "state": loc.get("state") or "",
                                                "country": loc.get("country") or "",
                                            }
                                        )

                                synthesized_for_ba = gen_data.get("_synthesized", {})
                                enriched_for_ba = gen_data.get("_enriched", {})
                                merged_for_ba = {}
                                if isinstance(enriched_for_ba, dict):
                                    merged_for_ba.update(enriched_for_ba)
                                if isinstance(synthesized_for_ba, dict):
                                    merged_for_ba.update(synthesized_for_ba)
                                budget_result = calculate_budget_allocation(
                                    total_budget=_bval_ba,
                                    roles=_roles_for_ba,
                                    locations=_locs_for_ba,
                                    industry=gen_data.get("industry", "General"),
                                    channel_percentages=channel_pcts,
                                    synthesized_data=merged_for_ba,
                                    knowledge_base=kb,
                                    collar_type=gen_data.get("_collar_type") or "",
                                    campaign_start_month=int(
                                        gen_data.get("campaign_start_month") or 0 or 0
                                    ),
                                )
                                gen_data["_budget_allocation"] = budget_result
                                logger.info(
                                    "Async budget allocation complete: %s",
                                    (
                                        list(budget_result.keys())
                                        if isinstance(budget_result, dict)
                                        else "N/A"
                                    ),
                                )
                            except Exception as ba_err:
                                logger.warning(
                                    "Async budget allocation failed (non-fatal): %s",
                                    ba_err,
                                )
                                gen_data["_budget_allocation"] = {}
                        else:
                            gen_data["_budget_allocation"] = {}

                        _step_timings["budget_allocation"] = round(
                            time.time() - _t_budget, 2
                        )
                        logger.info(
                            "PERF: Budget allocation took %.2fs",
                            _step_timings["budget_allocation"],
                        )

                        # ── Inject SlotOps LinkedIn Benchmarks (async path) ──
                        _t_slotops = time.time()
                        try:
                            from slotops_engine import get_linkedin_benchmarks_for_plan

                            _slotops_country_a = "United States"
                            _locs_for_slotops_a = gen_data.get("locations") or []
                            if _locs_for_slotops_a:
                                _first_loc_a = (
                                    _locs_for_slotops_a[0]
                                    if isinstance(_locs_for_slotops_a, list)
                                    else _locs_for_slotops_a
                                )
                                if isinstance(_first_loc_a, dict):
                                    _slotops_country_a = (
                                        _first_loc_a.get("country") or "United States"
                                    )
                                elif isinstance(_first_loc_a, str):
                                    _parts_a = [
                                        p.strip() for p in _first_loc_a.split(",")
                                    ]
                                    if len(_parts_a) >= 2:
                                        _slotops_country_a = _parts_a[-1]
                            _li_bench_a = get_linkedin_benchmarks_for_plan(
                                _slotops_country_a
                            )
                            if _li_bench_a:
                                gen_data["_slotops_linkedin_benchmarks"] = _li_bench_a
                                logger.info(
                                    "Async SlotOps LinkedIn benchmarks injected for %s",
                                    _slotops_country_a,
                                )
                        except ImportError:
                            logger.debug(
                                "slotops_engine not available -- skipping LinkedIn benchmarks (async)"
                            )
                        except Exception as _slotops_err_a:
                            logger.warning(
                                "Async SlotOps LinkedIn benchmarks failed (non-fatal): %s",
                                _slotops_err_a,
                            )

                        _step_timings["slotops_inject"] = round(
                            time.time() - _t_slotops, 2
                        )
                        logger.info(
                            "PERF: SlotOps inject took %.2fs",
                            _step_timings["slotops_inject"],
                        )

                        # ── Creative Quality Score (P1-16, async path) ──
                        try:
                            from creative_quality_score import score_creative_quality

                            gen_data["_creative_quality_score"] = (
                                score_creative_quality(gen_data)
                            )
                            logger.info(
                                "Async creative quality score: %s (%s)",
                                gen_data["_creative_quality_score"].get("score"),
                                gen_data["_creative_quality_score"].get("grade"),
                            )
                        except ImportError:
                            logger.debug(
                                "creative_quality_score not available -- skipping (async)"
                            )
                        except Exception as _cqs_err_a:
                            logger.warning(
                                "Async creative quality score failed (non-fatal): %s",
                                _cqs_err_a,
                            )

                        if time.time() > _gen_deadline:
                            raise TimeoutError(
                                "Plan generation exceeded 2-minute time limit"
                            )

                        with _generation_jobs_lock:
                            if jid in _generation_jobs:
                                _generation_jobs[jid]["progress_pct"] = 70
                                _generation_jobs[jid][
                                    "status_message"
                                ] = "Verifying plan data..."

                        # ── Gold Standard Quality Gates (async path) ──
                        # Wrapped in ThreadPoolExecutor with 30s timeout to
                        # prevent hanging at 70% "Verifying plan data..."
                        _t_gs = time.time()
                        try:
                            from gold_standard import apply_all_quality_gates

                            _gs_pool = ThreadPoolExecutor(max_workers=1)
                            _gs_future = _gs_pool.submit(
                                apply_all_quality_gates, gen_data
                            )
                            try:
                                _gs_future.result(timeout=_gs_timeout)
                            except TimeoutError:
                                logger.warning(
                                    "Async Gold Standard gates timed out after %ds -- continuing with partial data",
                                    _gs_timeout,
                                )
                                _gs_future.cancel()
                                gen_data["_gold_standard"] = (
                                    gen_data.get("_gold_standard") or {}
                                )
                            finally:
                                _gs_pool.shutdown(wait=False)
                        except ImportError:
                            logger.warning(
                                "gold_standard module not available -- skipping quality gates"
                            )
                            gen_data["_gold_standard"] = {}
                        except Exception as _gs_err:
                            logger.error(
                                "Async Gold Standard gates failed: %s",
                                _gs_err,
                                exc_info=True,
                            )
                            gen_data["_gold_standard"] = {}
                        _step_timings["gold_standard"] = round(time.time() - _t_gs, 2)
                        logger.info(
                            "PERF: Gold Standard gates took %.2fs",
                            _step_timings["gold_standard"],
                        )

                        # ── Deadline check BETWEEN quality gates and verification ──
                        if time.time() > _gen_deadline:
                            logger.warning(
                                "Deadline exceeded after quality gates -- skipping verification"
                            )
                            gen_data["_verification"] = {
                                "status": "skipped",
                                "reason": "deadline_exceeded_before_verification",
                            }
                        else:
                            # Gemini/LLM verification of key plan data points
                            # Also wrapped with 30s timeout to prevent hangs
                            _t_verify = time.time()
                            try:
                                _vf_pool = ThreadPoolExecutor(max_workers=1)
                                _vf_future = _vf_pool.submit(
                                    _verify_plan_data, gen_data
                                )
                                try:
                                    gen_data["_verification"] = _vf_future.result(
                                        timeout=_verify_timeout
                                    )
                                except TimeoutError:
                                    logger.warning(
                                        "Plan data verification timed out after 30s -- skipping"
                                    )
                                    _vf_future.cancel()
                                    gen_data["_verification"] = {
                                        "status": "skipped",
                                        "reason": "verification_timeout",
                                    }
                                finally:
                                    _vf_pool.shutdown(wait=False)
                            except Exception:
                                gen_data["_verification"] = {
                                    "status": "skipped",
                                    "reason": "verification_error",
                                }
                            _step_timings["verification"] = round(
                                time.time() - _t_verify, 2
                            )
                            logger.info(
                                "PERF: Verification took %.2fs",
                                _step_timings["verification"],
                            )

                        if time.time() > _gen_deadline:
                            # S27: Instead of raising, set flag and continue with truncated output
                            gen_data["_deadline_exceeded"] = True
                            logger.warning(
                                "Deadline exceeded before Excel for job %s -- continuing with truncated output",
                                jid,
                            )

                        # ── Inject international benchmarks for non-US plans ──
                        _target_region_for_intl = (
                            gen_data.get("target_region") or "us_only"
                        )
                        if _target_region_for_intl in (
                            "global",
                            "emea",
                            "apac",
                            "custom",
                        ):
                            try:
                                _kb_intl = load_knowledge_base().get(
                                    "international_benchmarks", {}
                                )
                                if _kb_intl:
                                    _intl_countries = _kb_intl.get("countries", {})
                                    _intl_regions = _kb_intl.get("regions", {})
                                    # Filter to relevant countries for the region
                                    _relevant_countries: dict = {}
                                    if _target_region_for_intl == "emea":
                                        _region_keys = ["emea"]
                                    elif _target_region_for_intl == "apac":
                                        _region_keys = ["apac"]
                                    elif _target_region_for_intl == "custom":
                                        _region_keys = list(_intl_regions.keys())
                                    else:  # global
                                        _region_keys = list(_intl_regions.keys())
                                    for _rk in _region_keys:
                                        _rcountries = _intl_regions.get(_rk, {}).get(
                                            "countries", []
                                        )
                                        for _ck in _rcountries:
                                            if _ck in _intl_countries:
                                                _relevant_countries[_ck] = (
                                                    _intl_countries[_ck]
                                                )
                                    gen_data["_intl_benchmarks"] = {
                                        "countries": _relevant_countries,
                                        "regions": {
                                            k: v
                                            for k, v in _intl_regions.items()
                                            if k in _region_keys
                                        },
                                        "source": "international_benchmarks_2026 (38 countries)",
                                    }
                                    logger.info(
                                        "Injected intl benchmarks: %d countries for region=%s",
                                        len(_relevant_countries),
                                        _target_region_for_intl,
                                    )
                            except Exception as _intl_err:
                                logger.error(
                                    "Failed to inject intl benchmarks: %s",
                                    _intl_err,
                                    exc_info=True,
                                )

                        # ── Cross-validation post-processing (async path) ──
                        # Detects internal inconsistencies in the enriched plan
                        # and auto-corrects where safe.  Non-fatal: failures
                        # here never block Excel/PPT generation.
                        _t_validator = time.time()
                        try:
                            from plan_validator import validate_plan

                            _val_result = validate_plan(gen_data)
                            logger.info(
                                "Plan validation (async): %d finding(s), %d auto-corrected",
                                _val_result.get("total_findings", 0),
                                _val_result.get("auto_corrections", 0),
                            )
                        except ImportError:
                            logger.debug("plan_validator not available -- skipping")
                        except Exception as _val_err:
                            logger.warning(
                                "Plan validation failed (async, non-fatal): %s",
                                _val_err,
                                exc_info=True,
                            )
                            gen_data["_validation"] = {
                                "findings": [],
                                "checks_run": 0,
                                "checks_failed": 1,
                                "error": str(_val_err),
                            }
                        _step_timings["plan_validator"] = round(
                            time.time() - _t_validator, 2
                        )
                        logger.info(
                            "PERF: Plan validation took %.2fs",
                            _step_timings["plan_validator"],
                        )

                        # Excel generation -- wrapped in 60s timeout to prevent
                        # indefinite hangs on LLM calls during cold starts
                        _t_excel = time.time()
                        _excel_pool = ThreadPoolExecutor(
                            max_workers=1, thread_name_prefix="async-excel"
                        )
                        try:
                            if generate_excel_v2 is not None:
                                _excel_future = _excel_pool.submit(
                                    generate_excel_v2,
                                    gen_data,
                                    research_mod=research,
                                    load_kb_fn=load_knowledge_base,
                                    classify_tier_fn=classify_role_tier,
                                    fetch_logo_fn=fetch_client_logo,
                                )
                            elif generate_excel is not None:
                                _excel_future = _excel_pool.submit(
                                    generate_excel, gen_data
                                )
                            else:
                                raise RuntimeError("No Excel generator available")
                            excel_bytes = _excel_future.result(timeout=_excel_timeout)
                            logger.info(
                                "Async Excel generated (%d bytes)",
                                len(excel_bytes) if excel_bytes else 0,
                            )
                        except TimeoutError:
                            logger.error(
                                "Excel generation timed out after %ds for job %s",
                                _excel_timeout,
                                jid,
                            )
                            _excel_future.cancel()
                            raise TimeoutError(
                                "Excel generation exceeded 60-second limit"
                            )
                        finally:
                            _excel_pool.shutdown(wait=False, cancel_futures=True)

                        _step_timings["excel"] = round(time.time() - _t_excel, 2)
                        logger.info(
                            "PERF: Excel generation took %.2fs", _step_timings["excel"]
                        )

                        with _generation_jobs_lock:
                            if jid in _generation_jobs:
                                _generation_jobs[jid]["progress_pct"] = 80
                                _generation_jobs[jid][
                                    "status_message"
                                ] = "Generating Excel report..."

                        if time.time() > _gen_deadline:
                            raise TimeoutError(
                                "Plan generation exceeded 2-minute time limit"
                            )

                        # PPT generation -- skip if caller only wants Excel
                        # Accepted output_format values: "excel", "ppt", "all", "google" (default: "all")
                        # S49 P2-8: "google" generates Excel+PPT bundle (same as "all"),
                        # then additionally attempts async Google Sheets export. The frontend
                        # checks /api/export/status to enable/disable the Google toggle --
                        # button is greyed out when GOOGLE_SHEETS_CREDENTIALS / GOOGLE_SLIDES_CREDENTIALS_B64
                        # is not configured.
                        _requested_format = (
                            (gen_data.get("output_format") or "all").lower().strip()
                        )
                        _wants_google_export = _requested_format == "google"
                        if _requested_format == "google":
                            _requested_format = "all"
                        client_name = re.sub(
                            r"_+",
                            "_",  # S27: Collapse multiple underscores
                            re.sub(
                                r"[^a-zA-Z0-9_\-]",
                                "_",
                                gen_data.get("client_name") or "Client",
                            ),
                        ).strip("_")
                        pptx_bytes = None
                        if generate_pptx is not None and _requested_format in (
                            "all",
                            "ppt",
                            "bundle",
                        ):
                            with _generation_jobs_lock:
                                if jid in _generation_jobs:
                                    _generation_jobs[jid]["progress_pct"] = 90
                                    _generation_jobs[jid][
                                        "status_message"
                                    ] = "Creating strategy deck..."
                            # PPT wrapped in adaptive timeout
                            _ppt_timeout = 90  # S29 v2: generous -- quality PPT
                            _ppt_pool = ThreadPoolExecutor(
                                max_workers=1, thread_name_prefix="async-ppt"
                            )
                            try:
                                # S48 FIX: Route through DeckGenerator for Joveo template + Google Slides
                                # Previously called generate_pptx directly, bypassing Tier 1 (Google Slides)
                                if _deck_generator is not None:

                                    def _async_deck_gen(d):
                                        """Wrapper: DeckGenerator returns (bytes, provider); we need just bytes."""
                                        pptx_b, provider = _deck_generator.generate(d)
                                        logger.info(
                                            "Async PPT generated via %s (%d bytes)",
                                            provider,
                                            len(pptx_b),
                                        )
                                        return pptx_b

                                    _ppt_future = _ppt_pool.submit(
                                        _async_deck_gen, gen_data
                                    )
                                else:
                                    _ppt_future = _ppt_pool.submit(
                                        generate_pptx, gen_data
                                    )
                                pptx_bytes = _ppt_future.result(timeout=_ppt_timeout)
                            except TimeoutError:
                                logger.warning(
                                    "PPT generation timed out after %ds for job %s -- skipping PPT",
                                    _ppt_timeout,
                                    jid,
                                )
                                _ppt_future.cancel()
                                pptx_bytes = None
                            except Exception:
                                pptx_bytes = None
                            finally:
                                _ppt_pool.shutdown(wait=False, cancel_futures=True)
                        elif _requested_format == "excel":
                            logger.info("Skipping PPT generation (output_format=excel)")

                        # PERF: Log total pipeline timing summary
                        _pipeline_total = round(time.time() - _pipeline_t0, 2)
                        _step_timings["total"] = _pipeline_total
                        logger.info(
                            "PERF: Pipeline complete in %.2fs | %s",
                            _pipeline_total,
                            " | ".join(
                                f"{k}={v:.2f}s" for k, v in _step_timings.items()
                            ),
                        )

                        if pptx_bytes:
                            zip_buffer = io.BytesIO()
                            with zipfile.ZipFile(
                                zip_buffer, "w", zipfile.ZIP_DEFLATED
                            ) as zf:
                                zf.writestr(
                                    f"{client_name}_Media_Plan.xlsx", excel_bytes
                                )
                                zf.writestr(
                                    f"{client_name}_Strategy_Deck.pptx", pptx_bytes
                                )
                            result_bytes = zip_buffer.getvalue()
                            result_ct = "application/zip"
                            result_fn = f"{client_name}_Media_Plan_Bundle.zip"
                        else:
                            result_bytes = excel_bytes
                            result_ct = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                            result_fn = f"{client_name}_Media_Plan.xlsx"

                        with _generation_jobs_lock:
                            if jid in _generation_jobs:
                                _generation_jobs[jid].update(
                                    {
                                        "status": "completed",
                                        "progress_pct": 100,
                                        "status_message": "Complete",
                                        "result_bytes": result_bytes,
                                        "result_content_type": result_ct,
                                        "result_filename": result_fn,
                                    }
                                )
                        logger.info(
                            "Async job %s completed (%d bytes)", jid, len(result_bytes)
                        )

                        # ── S47: Persist ZIP to Supabase for restart-safe Slack links ──
                        try:
                            from supabase_client import get_client as _sb_client
                            import base64 as _b64_persist

                            _sb = _sb_client()
                            if _sb and result_bytes:
                                _zip_b64 = _b64_persist.b64encode(result_bytes).decode()
                                _sb.table("nova_generated_plans").upsert(
                                    {
                                        "job_id": jid,
                                        "zip_data": _zip_b64,
                                        "filename": result_fn,
                                        "client_name": gen_data.get("client_name")
                                        or "",
                                        "industry": (
                                            gen_data.get("industry_label")
                                            or gen_data.get("industry")
                                            or ""
                                        ),
                                        "user_email": gen_data.get("requester_email")
                                        or "",
                                        "created_at": time.strftime(
                                            "%Y-%m-%dT%H:%M:%SZ", time.gmtime()
                                        ),
                                    }
                                ).execute()
                                logger.info(
                                    "S47: Persisted plan ZIP to Supabase for job %s",
                                    jid,
                                )
                        except Exception as _sb_err:
                            logger.debug(
                                "S47: Failed to persist plan to Supabase (non-critical): %s",
                                _sb_err,
                            )

                        # ── Upload ZIP to Google Drive for permanent Slack link ──
                        _drive_download_url = ""
                        try:
                            _drive_download_url = _upload_plan_to_drive(
                                result_bytes, result_fn, result_ct
                            )
                            if _drive_download_url:
                                with _generation_jobs_lock:
                                    if jid in _generation_jobs:
                                        _generation_jobs[jid][
                                            "drive_url"
                                        ] = _drive_download_url
                                logger.info(
                                    "Plan uploaded to Drive for job %s: %s",
                                    jid,
                                    _drive_download_url,
                                )
                        except Exception as drive_exc:
                            logger.warning(
                                "Drive upload failed for job %s (non-critical): %s",
                                jid,
                                drive_exc,
                            )

                        # ── Slack notification: plan generated ──
                        try:
                            from slack_plan_notifier import notify_plan_generated

                            _budget_alloc = gen_data.get("_budget_allocation") or {}
                            _ch_allocs = _budget_alloc.get("channel_allocations") or {}
                            _channels_list = (
                                list(_ch_allocs.keys()) if _ch_allocs else []
                            )
                            _roles = (
                                gen_data.get("target_roles")
                                or gen_data.get("roles")
                                or []
                            )
                            _locs = gen_data.get("locations") or []
                            _gen_elapsed = time.time() - gen_data.get(
                                "_gen_created", time.time()
                            )

                            notify_plan_generated(
                                {
                                    "client_name": gen_data.get("client_name")
                                    or "Unknown",
                                    "industry": (
                                        gen_data.get("industry_label")
                                        or gen_data.get("industry")
                                        or "N/A"
                                    ),
                                    "budget": gen_data.get("budget")
                                    or gen_data.get("budget_range")
                                    or "N/A",
                                    "user_name": gen_data.get("requester_name")
                                    or "Unknown",
                                    "user_email": gen_data.get("requester_email")
                                    or "Unknown",
                                    "role_count": (
                                        len(_roles) if isinstance(_roles, list) else 0
                                    ),
                                    "location_count": (
                                        len(_locs) if isinstance(_locs, list) else 0
                                    ),
                                    "channel_count": len(_channels_list),
                                    "channels_list": _channels_list,
                                    "generation_time_seconds": _gen_elapsed,
                                    "job_id": jid,
                                    "drive_url": _drive_download_url,
                                    "filename": result_fn,
                                    "timestamp": time.strftime(
                                        "%Y-%m-%d %H:%M:%S UTC", time.gmtime()
                                    ),
                                }
                            )
                        except Exception:
                            pass  # Slack notification is best-effort

                    except Exception as async_err:
                        logger.error(
                            "Async job %s failed: %s", jid, async_err, exc_info=True
                        )
                        with _generation_jobs_lock:
                            if jid in _generation_jobs:
                                _generation_jobs[jid].update(
                                    {
                                        "status": "failed",
                                        "error": str(async_err),
                                        "result_bytes": None,
                                    }
                                )
                        # Send email alert for async generation failure (matches sync path)
                        try:
                            from email_alerts import send_generation_failure_alert

                            _client = (
                                gen_data.get("client_name", "Unknown")
                                if gen_data
                                else "Unknown"
                            )
                            send_generation_failure_alert(_client, str(async_err), jid)
                        except Exception:
                            pass  # email alerts are best-effort

                t = threading.Thread(
                    target=_async_generate,
                    args=(job_id, data, request_id),
                    daemon=True,
                    name=f"async-gen-{job_id}",
                )
                t.start()
                self._send_json(
                    {
                        "job_id": job_id,
                        "status": "processing",
                        "poll_url": f"/api/jobs/{job_id}",
                    }
                )
                # ── PostHog: Track async media plan generation ──
                self._ph_track(
                    "media_plan_generated", {"async": True, "job_id": job_id}
                )
                return

            # ── Process uploaded files (briefs, transcripts, historical data) ──
            try:
                from file_processor import (
                    extract_all_texts,
                    parse_historical_data,
                    summarize_transcript_context,
                )

                # 1. Brief file uploads -> append extracted text to use_case
                _upload_counts = {}  # Track counts before popping for PostHog
                uploaded_briefs = data.pop("uploaded_briefs", []) or []
                _upload_counts["briefs"] = len(uploaded_briefs)
                if uploaded_briefs:
                    brief_text = extract_all_texts(uploaded_briefs)
                    if brief_text:
                        existing_use_case = data.get("use_case") or ""
                        data["use_case"] = (
                            existing_use_case
                            + "\n\n--- UPLOADED BRIEF DOCUMENTS ---\n"
                            + brief_text
                        ).strip()
                        logger.info(
                            "Appended %d chars from %d uploaded brief file(s)",
                            len(brief_text),
                            len(uploaded_briefs),
                        )

                # 2. Call transcript -> extract context and append to use_case
                call_transcript = data.pop("call_transcript", "") or ""
                uploaded_transcripts = data.pop("uploaded_transcripts", []) or []
                _upload_counts["transcripts"] = len(uploaded_transcripts)
                transcript_text = call_transcript
                if uploaded_transcripts:
                    file_transcript = extract_all_texts(uploaded_transcripts)
                    transcript_text = (
                        (transcript_text + "\n\n" + file_transcript).strip()
                        if transcript_text
                        else file_transcript
                    )
                if transcript_text and len(transcript_text.strip()) > 30:
                    context_summary = summarize_transcript_context(transcript_text)
                    if context_summary:
                        existing_use_case = data.get("use_case") or ""
                        data["use_case"] = (
                            existing_use_case
                            + "\n\n--- SALES CALL CONTEXT ---\n"
                            + context_summary
                        ).strip()
                        logger.info(
                            "Added %d chars of call transcript context",
                            len(context_summary),
                        )

                # 3. Historical performance data -> parse and store for downstream use
                uploaded_historical = data.pop("uploaded_historical", []) or []
                _upload_counts["historical"] = len(uploaded_historical)
                if uploaded_historical:
                    historical_data = parse_historical_data(uploaded_historical)
                    if historical_data:
                        data["historical_performance"] = historical_data
                        logger.info(
                            "Parsed historical data: %d platforms, $%.0f total spend",
                            historical_data.get("platform_count") or 0,
                            historical_data.get("total_spend") or 0,
                        )
            except ImportError:
                logger.warning(
                    "file_processor module not available -- skipping file uploads"
                )
                # Remove file fields to avoid bloating downstream processing
                data.pop("uploaded_briefs", None)
                data.pop("uploaded_transcripts", None)
                data.pop("uploaded_historical", None)
                data.pop("call_transcript", None)
            except Exception as _fp_err:
                logger.warning("File processing error (non-fatal): %s", _fp_err)
                data.pop("uploaded_briefs", None)
                data.pop("uploaded_transcripts", None)
                data.pop("uploaded_historical", None)
                data.pop("call_transcript", None)

            # ── P0 Fix: Capture timeline, hire_volume, and notes from input ──
            # campaign_duration: normalize from frontend dropdown or API input
            raw_duration = str(
                data.get("campaign_duration") or data.get("timeline") or ""
            ).strip()
            if raw_duration:
                data["campaign_duration"] = raw_duration
            else:
                data["campaign_duration"] = "Not specified"

            # campaign_start_month: int 1-12, 0 or absent = current month
            raw_csm = data.get("campaign_start_month") or 0
            try:
                campaign_start_month = int(raw_csm) if raw_csm else 0
            except (ValueError, TypeError):
                campaign_start_month = 0
            if campaign_start_month < 0 or campaign_start_month > 12:
                campaign_start_month = 0
            if campaign_start_month == 0:
                import datetime as _dt_csm

                campaign_start_month = _dt_csm.datetime.now().month
            data["campaign_start_month"] = campaign_start_month

            # hire_volume: capture from form field or parse from notes
            raw_hire_vol = str(data.get("hire_volume") or "").strip()
            raw_notes = (data.get("notes") or "").strip()

            if raw_hire_vol:
                data["hire_volume"] = raw_hire_vol
            elif raw_notes:
                # Parse hire volume from notes (e.g. "5000+ hires", "10,000 positions", "500 teachers")
                # Keyword group is REQUIRED to avoid matching budgets ($50,000), dates (2026), etc.
                hv_match = re.search(
                    r"(\d[\d,]*)\+?\s*(?:hires?|positions?|roles?|openings?|teachers?|nurses?|drivers?|employees?|workers?|associates?|people|staff|headcount)",
                    raw_notes,
                    re.IGNORECASE,
                )
                if hv_match:
                    hv_num = hv_match.group(1).replace(",", "")
                    try:
                        hv_int = int(hv_num)
                        if hv_int > 1_000_000:
                            hv_int = 1_000_000  # Cap at 1M
                        if hv_int >= 5:  # Only treat as hire volume if >= 5
                            # Preserve the original text snippet for context
                            hv_text = hv_match.group(0).strip()
                            if "+" in raw_notes[hv_match.start() : hv_match.end() + 2]:
                                hv_text = f"{hv_int:,}+"
                            else:
                                hv_text = f"{hv_int:,}"
                            data["hire_volume"] = f"{hv_text} hires"
                    except ValueError:
                        pass
            if not data.get("hire_volume"):
                data["hire_volume"] = "Not specified"

            # Store notes for downstream use
            if raw_notes:
                data["notes"] = raw_notes

            # ── Normalize roles early: dict-of-dicts -> list-of-strings ──
            # API/form can send roles as [{"title":"...","location":"..."}] or ["string"]
            for _rkey in ("roles", "target_roles"):
                _rlist = data.get(_rkey) or []
                if isinstance(_rlist, list) and _rlist and isinstance(_rlist[0], dict):
                    data[_rkey] = [
                        (r.get("title") or r.get("role") or str(r)) for r in _rlist
                    ]

            # Compute campaign_weeks from campaign_duration for timeline phasing
            duration_str = str(data.get("campaign_duration") or "" or "")
            campaign_weeks = 12  # default
            dur_lower = duration_str.lower()
            # Order matters: check more specific/longer ranges before shorter ones
            if (
                "2-5 year" in dur_lower
                or "long-term" in dur_lower
                or "long term" in dur_lower
            ):
                campaign_weeks = 156
            elif "1-2 year" in dur_lower or "2 year" in dur_lower:
                campaign_weeks = 80
            elif (
                "6-12 month" in dur_lower
                or "9 month" in dur_lower
                or "12 month" in dur_lower
                or "1 year" in dur_lower
            ):
                campaign_weeks = 48
            elif (
                "3-6 month" in dur_lower
                or "4 month" in dur_lower
                or "5 month" in dur_lower
                or "6 month" in dur_lower
            ):
                campaign_weeks = 24
            elif (
                "1-3 month" in dur_lower
                or "1 month" in dur_lower
                or "2 month" in dur_lower
                or "3 month" in dur_lower
            ):
                campaign_weeks = 12
            elif "ongoing" in dur_lower:
                campaign_weeks = 52  # annual cycle
            else:
                # Try to parse weeks directly (e.g. "26 weeks", "16 weeks")
                wk_match = re.search(r"(\d+)\s*week", dur_lower)
                if wk_match:
                    campaign_weeks = int(wk_match.group(1))
                else:
                    # Try months (e.g. "6 months")
                    mo_match = re.search(r"(\d+)\s*month", dur_lower)
                    if mo_match:
                        campaign_weeks = int(mo_match.group(1)) * 4
            data["campaign_weeks"] = campaign_weeks

            # ── Phase 0: Canonical Taxonomy Normalization ──
            # Run the standardizer on all input fields BEFORE they enter
            # the enrichment/synthesis pipeline.  This ensures every
            # downstream lookup (KB, channels_db, BLS, etc.) uses
            # consistent canonical keys instead of variant strings.
            if _STANDARDIZER_AVAILABLE:
                try:
                    # -- Normalize industry --
                    raw_ind = data.get("industry") or ""
                    canonical_ind = std_normalize_industry(raw_ind)
                    data["_industry_original"] = raw_ind
                    data["_industry_canonical"] = canonical_ind
                    ind_meta = CANONICAL_INDUSTRIES.get(canonical_ind, {})
                    data["_industry_legacy_key"] = ind_meta.get(
                        "deep_bench_key", ind_meta.get("aliases", [""])[0]
                    )
                    data["_industry_kb_key"] = ind_meta.get("kb_key") or ""

                    # Update the primary industry key so the enrichment
                    # pipeline (api_enrichment, data_synthesizer, etc.)
                    # receives a canonical key instead of the raw variant.
                    # The classify_industry call downstream will further
                    # refine this to the legacy_key for backward compat.
                    if canonical_ind and canonical_ind != "general":
                        data["industry"] = canonical_ind
                    elif data.get("_industry_legacy_key"):
                        data["industry"] = data["_industry_legacy_key"]

                    # -- Normalize roles (attach SOC + tier metadata) --
                    for role_key in ("target_roles", "roles"):
                        raw_roles = data.get(role_key) or []
                        if isinstance(raw_roles, list):
                            for r in raw_roles:
                                if isinstance(r, dict):
                                    title = r.get("title") or ""
                                    r["_canonical_role"] = std_normalize_role(title)
                                    r["_soc_code"] = std_get_soc_code(title)
                                    r["_role_tier"] = std_get_role_tier(title)

                    # -- Normalize locations (attach region/market keys) --
                    raw_locs = data.get("locations") or []
                    if isinstance(raw_locs, list):
                        parsed_locs = []
                        for loc in raw_locs:
                            loc_str = loc if isinstance(loc, str) else ""
                            if isinstance(loc, dict):
                                loc_str = ", ".join(
                                    filter(
                                        None,
                                        [
                                            loc.get("city") or "",
                                            loc.get("state") or "",
                                            loc.get("country") or "",
                                        ],
                                    )
                                )
                            parsed_locs.append(std_normalize_location(loc_str))
                        data["_locations_parsed"] = parsed_locs

                    logger.info(
                        "Standardizer: industry=%s -> canonical=%s (kb_key=%s)",
                        raw_ind,
                        canonical_ind,
                        data.get("_industry_kb_key") or "",
                    )
                except Exception as e:
                    logger.warning(
                        "Standardizer normalization failed (non-fatal): %s", e
                    )

            # ── Parallel API Enrichment for /api/generate ──
            # All enrichment sources run concurrently with a 20s global timeout.
            # Previously serial (worst-case 60s+), now parallel (worst-case ~20s).
            _span_fn = getattr(self, "_sentry_span", lambda o, n: _nullctx())
            enriched: dict = {}
            _market_ctx: dict = {}

            # -- Extract shared params once for all enrichment tasks --
            _roles_for_api = data.get("target_roles") or data.get("roles") or []

            # Build normalized list of roles (up to 5) with title + SOC code
            _roles_normalized: list[dict[str, str]] = []
            if isinstance(_roles_for_api, list):
                for _r in _roles_for_api[:5]:
                    if isinstance(_r, dict):
                        _t = str(_r.get("title") or "")
                        _s = str(_r.get("_soc_code") or "")
                        if _t:
                            _roles_normalized.append({"title": _t, "soc": _s})
                    elif isinstance(_r, str) and _r.strip():
                        _roles_normalized.append({"title": _r.strip(), "soc": ""})
            elif _roles_for_api:
                _roles_normalized.append({"title": str(_roles_for_api), "soc": ""})

            # Backward-compat: first role string for callables that need a single value
            _role_str = _roles_normalized[0]["title"] if _roles_normalized else ""

            _locs_for_api = data.get("locations") or []
            _state_code = ""
            if _locs_for_api:
                _first_loc = (
                    _locs_for_api[0]
                    if isinstance(_locs_for_api, list)
                    else str(_locs_for_api)
                )
                if isinstance(_first_loc, dict):
                    _state_code = _first_loc.get("state") or ""
                elif isinstance(_first_loc, str):
                    _parts = [p.strip() for p in _first_loc.split(",")]
                    _state_code = _parts[1] if len(_parts) > 1 else ""
                if len(_state_code) > 2:
                    _state_code = ""  # Only 2-letter state codes

            _soc_code = ""
            for _rr in (_roles_for_api if isinstance(_roles_for_api, list) else []):
                if isinstance(_rr, dict):
                    _soc_code = _rr.get("_soc_code") or ""
                    break

            # JobSpy location extraction
            _js_loc = "USA"
            if isinstance(_locs_for_api, list) and _locs_for_api:
                _js_first_loc = _locs_for_api[0]
                if isinstance(_js_first_loc, str):
                    _js_loc = _js_first_loc
                elif isinstance(_js_first_loc, dict):
                    _js_loc = (
                        _js_first_loc.get("city") or _js_first_loc.get("state") or "USA"
                    )

            # -- Define individual enrichment callables --
            # Each returns (name, result_dict_or_none) and isolates its own errors.

            def _call_enrich_data() -> "tuple[str, Optional[dict]]":
                """Run the api_enrichment.enrich_data module."""
                if enrich_data is None:
                    return ("api_enrichment", None)
                try:
                    _rid = getattr(self, "_request_id", None)
                    result = (
                        enrich_data(data, request_id=_rid)
                        if _rid
                        else enrich_data(data)
                    ) or {}
                    logger.info(
                        "API enrichment complete: %s",
                        result.get("enrichment_summary", {}),
                    )
                    return ("api_enrichment", result)
                except Exception as exc:
                    logger.warning("API enrichment failed (non-fatal): %s", exc)
                    return ("api_enrichment", None)

            def _call_adzuna() -> "tuple[str, Optional[dict]]":
                """Fetch Adzuna job count for ALL roles (up to 5)."""
                if not (
                    _api_integrations_available and _api_adzuna and _roles_normalized
                ):
                    return ("adzuna_job_count", None)
                per_role: list[dict[str, Any]] = []
                for _ri in _roles_normalized:
                    try:
                        count = _api_adzuna.get_job_count(_ri["title"], "us")
                        if count:
                            per_role.append({"role": _ri["title"], "count": count})
                    except (
                        urllib.error.URLError,
                        OSError,
                        ValueError,
                        TypeError,
                    ) as exc:
                        logger.error(
                            f"Adzuna enrichment for role '{_ri['title']}' failed: {exc}",
                            exc_info=True,
                        )
                if not per_role:
                    return ("adzuna_job_count", None)
                total = sum(
                    r["count"] for r in per_role if isinstance(r["count"], (int, float))
                )
                logger.info(
                    f"Enriched /api/generate with adzuna data for {len(per_role)} role(s)"
                )
                return (
                    "adzuna_job_count",
                    {
                        "total": total,
                        "per_role": per_role,
                    },
                )

            def _call_bls() -> "tuple[str, Optional[dict]]":
                """Fetch BLS occupational employment data for ALL roles with SOC codes."""
                if not (_api_integrations_available and _api_bls):
                    return ("bls_occupational_employment", None)
                per_role: list[dict[str, Any]] = []
                for _ri in _roles_normalized:
                    _soc = _ri.get("soc") or ""
                    if not _soc:
                        continue
                    try:
                        oes = _api_bls.get_occupational_employment(_soc)
                        if oes:
                            per_role.append(
                                {"role": _ri["title"], "soc": _soc, "data": oes}
                            )
                    except (
                        urllib.error.URLError,
                        OSError,
                        ValueError,
                        TypeError,
                    ) as exc:
                        logger.error(
                            f"BLS enrichment for role '{_ri['title']}' (SOC {_soc}) failed: {exc}",
                            exc_info=True,
                        )
                if not per_role:
                    return ("bls_occupational_employment", None)
                logger.info(
                    f"Enriched /api/generate with BLS OES data for {len(per_role)} role(s)"
                )
                # Return first role's data at top level for backward compat, plus per_role array
                return (
                    "bls_occupational_employment",
                    {
                        **per_role[0]["data"],
                        "per_role": per_role,
                    },
                )

            def _call_fred() -> "tuple[str, Optional[dict]]":
                """Fetch FRED unemployment rate (optionally by state)."""
                if not (_api_integrations_available and _api_fred):
                    return ("fred_unemployment", None)
                try:
                    unemp = _api_fred.get_unemployment_rate(
                        state_code=_state_code if _state_code else None
                    )
                    if unemp:
                        logger.info(
                            "Enriched /api/generate with fred unemployment data"
                        )
                        return ("fred_unemployment", unemp)
                    return ("fred_unemployment", None)
                except (urllib.error.URLError, OSError, ValueError, TypeError) as exc:
                    logger.error(
                        "FRED enrichment for /api/generate failed: %s",
                        exc,
                        exc_info=True,
                    )
                    return ("fred_unemployment", None)

            def _call_jobspy() -> "tuple[str, Optional[dict]]":
                """Fetch JobSpy market stats for ALL roles (posting volume + salary)."""
                if not (
                    _jobspy_available and _jobspy_market_stats and _roles_normalized
                ):
                    return ("jobspy_market_stats", None)
                per_role: list[dict[str, Any]] = []
                for _ri in _roles_normalized:
                    try:
                        stats = _jobspy_market_stats(_ri["title"], _js_loc)
                        if stats:
                            per_role.append({"role": _ri["title"], **stats})
                    except (ValueError, TypeError, KeyError, OSError) as exc:
                        logger.error(
                            f"JobSpy enrichment for role '{_ri['title']}' failed: {exc}",
                            exc_info=True,
                        )
                if not per_role:
                    return ("jobspy_market_stats", None)
                # Aggregate totals
                total_postings = sum(r.get("total_postings") or 0 for r in per_role)

                # Fallback: when APIs return 0 postings, estimate from
                # industry averages so downstream never shows "0 postings".
                if total_postings == 0:
                    _industry_lower = str(data.get("industry") or "").lower()
                    _POSTINGS_FALLBACK = {
                        "tech": 50000,
                        "technology": 50000,
                        "tech_engineering": 50000,
                        "healthcare": 80000,
                        "healthcare_medical": 80000,
                        "medical": 80000,
                        "finance": 30000,
                        "finance_banking": 30000,
                        "banking": 30000,
                    }
                    _fallback_per_role = 25000  # default
                    for _fb_key, _fb_val in _POSTINGS_FALLBACK.items():
                        if _fb_key in _industry_lower:
                            _fallback_per_role = _fb_val
                            break
                    total_postings = _fallback_per_role * len(per_role)
                    logger.info(
                        f"JobSpy returned 0 postings; using industry fallback "
                        f"({_fallback_per_role}/role x {len(per_role)} roles = {total_postings})"
                    )

                avg_salaries = [
                    r["avg_salary"] for r in per_role if r.get("avg_salary")
                ]
                avg_salary = (
                    round(sum(avg_salaries) / len(avg_salaries), 2)
                    if avg_salaries
                    else None
                )
                logger.info(
                    f"Enriched /api/generate with jobspy stats for {len(per_role)} role(s) "
                    f"({total_postings} postings, avg_salary={avg_salary})"
                )
                return (
                    "jobspy_market_stats",
                    {
                        "total_postings": total_postings,
                        "avg_salary": avg_salary,
                        "per_role": per_role,
                    },
                )

            # ── Additional enrichment callables (S50: wire configured APIs) ──

            def _call_onet() -> "tuple[str, Optional[dict]]":
                """Fetch O*NET skills, knowledge, and occupation details for ALL roles."""
                if not (
                    _api_integrations_available and _api_onet and _roles_normalized
                ):
                    return ("onet_skills", None)
                per_role: list[dict[str, Any]] = []
                for _ri in _roles_normalized[:3]:  # Cap at 3 roles to stay fast
                    _soc = _ri.get("soc") or ""
                    _title = _ri.get("title") or ""
                    try:
                        # If no SOC code, try to look up by title
                        if not _soc:
                            occs = _api_onet.search_occupations(_title)
                            if occs and isinstance(occs, list) and occs:
                                _soc = (
                                    (occs[0].get("code") or "")
                                    if isinstance(occs[0], dict)
                                    else ""
                                )
                        if not _soc:
                            continue
                        skills = _api_onet.get_skills(_soc)
                        knowledge = _api_onet.get_knowledge(_soc)
                        details = _api_onet.get_occupation_details(_soc)
                        entry: dict[str, Any] = {"role": _title, "soc": _soc}
                        if skills:
                            entry["skills"] = skills[:10]  # Top 10 skills
                        if knowledge:
                            entry["knowledge"] = knowledge[:10]
                        if details and isinstance(details, dict):
                            entry["title_onet"] = details.get("title") or _title
                            entry["description"] = (details.get("description") or "")[
                                :500
                            ]
                        if len(entry) > 2:  # Has data beyond role+soc
                            per_role.append(entry)
                    except (
                        urllib.error.URLError,
                        OSError,
                        ValueError,
                        TypeError,
                    ) as exc:
                        logger.error(
                            f"O*NET enrichment for role '{_title}' (SOC {_soc}) failed: {exc}",
                            exc_info=True,
                        )
                if not per_role:
                    return ("onet_skills", None)
                logger.info(
                    f"Enriched /api/generate with O*NET data for {len(per_role)} role(s)"
                )
                return ("onet_skills", {"per_role": per_role})

            def _call_jooble() -> "tuple[str, Optional[dict]]":
                """Fetch Jooble international job postings (skip for US-only plans)."""
                if not (
                    _api_integrations_available and _api_jooble and _roles_normalized
                ):
                    return ("jooble_international", None)
                # Only call Jooble for international locations
                _locs_lower = " ".join(
                    str(l).lower()
                    for l in (_locs_for_api if isinstance(_locs_for_api, list) else [])
                )
                _us_only_keywords = {"us", "usa", "united states", "america"}
                _loc_words = set(_locs_lower.replace(",", " ").split())
                # If ALL location words are US-only, skip Jooble
                if _loc_words and _loc_words.issubset(
                    _us_only_keywords
                    | {
                        "al",
                        "ak",
                        "az",
                        "ar",
                        "ca",
                        "co",
                        "ct",
                        "de",
                        "fl",
                        "ga",
                        "hi",
                        "id",
                        "il",
                        "in",
                        "ia",
                        "ks",
                        "ky",
                        "la",
                        "me",
                        "md",
                        "ma",
                        "mi",
                        "mn",
                        "ms",
                        "mo",
                        "mt",
                        "ne",
                        "nv",
                        "nh",
                        "nj",
                        "nm",
                        "ny",
                        "nc",
                        "nd",
                        "oh",
                        "ok",
                        "or",
                        "pa",
                        "ri",
                        "sc",
                        "sd",
                        "tn",
                        "tx",
                        "ut",
                        "vt",
                        "va",
                        "wa",
                        "wv",
                        "wi",
                        "wy",
                    }
                ):
                    return ("jooble_international", None)
                per_role: list[dict[str, Any]] = []
                for _ri in _roles_normalized[:3]:
                    try:
                        _loc_str = ""
                        if isinstance(_locs_for_api, list) and _locs_for_api:
                            _first = _locs_for_api[0]
                            _loc_str = (
                                _first if isinstance(_first, str) else str(_first)
                            )
                        result = _api_jooble.search_jobs(_ri["title"], _loc_str)
                        if result and isinstance(result, dict):
                            total = result.get("totalCount") or 0
                            per_role.append(
                                {"role": _ri["title"], "total_count": total}
                            )
                    except (
                        urllib.error.URLError,
                        OSError,
                        ValueError,
                        TypeError,
                    ) as exc:
                        logger.error(
                            f"Jooble enrichment for role '{_ri['title']}' failed: {exc}",
                            exc_info=True,
                        )
                if not per_role:
                    return ("jooble_international", None)
                total = sum(r.get("total_count") or 0 for r in per_role)
                logger.info(
                    f"Enriched /api/generate with Jooble data for {len(per_role)} role(s) ({total} intl postings)"
                )
                return ("jooble_international", {"total": total, "per_role": per_role})

            def _call_bea() -> "tuple[str, Optional[dict]]":
                """Fetch BEA GDP/income data for the campaign state."""
                if not (_api_integrations_available and _api_bea):
                    return ("bea_economics", None)
                try:
                    result: dict[str, Any] = {}
                    # State-level GDP
                    gdp = _api_bea.get_gdp_by_state_all(year="LAST5")
                    if gdp and isinstance(gdp, dict):
                        result["gdp_by_state"] = gdp
                    # Personal income by state
                    income = _api_bea.get_personal_income_all_states(year="LAST5")
                    if income and isinstance(income, dict):
                        result["personal_income"] = income
                    if not result:
                        return ("bea_economics", None)
                    logger.info("Enriched /api/generate with BEA economic data")
                    return ("bea_economics", result)
                except (urllib.error.URLError, OSError, ValueError, TypeError) as exc:
                    logger.error(f"BEA enrichment failed: {exc}", exc_info=True)
                    return ("bea_economics", None)

            def _call_census() -> "tuple[str, Optional[dict]]":
                """Fetch Census workforce demographics for the campaign state."""
                if not (_api_integrations_available and _api_census and _state_code):
                    return ("census_demographics", None)
                try:
                    from api_integrations import _resolve_state_fips

                    _fips = _resolve_state_fips(_state_code)
                    if not _fips:
                        return ("census_demographics", None)
                    result: dict[str, Any] = {}
                    profile = _api_census.get_workforce_profile(_state_code)
                    if profile and isinstance(profile, dict):
                        result["workforce_profile"] = profile
                    demographics = _api_census.get_demographics(_fips)
                    if demographics and isinstance(demographics, dict):
                        result["demographics"] = demographics
                    if not result:
                        return ("census_demographics", None)
                    logger.info(
                        f"Enriched /api/generate with Census data for state={_state_code}"
                    )
                    return ("census_demographics", result)
                except (
                    urllib.error.URLError,
                    OSError,
                    ValueError,
                    TypeError,
                    ImportError,
                ) as exc:
                    logger.error(f"Census enrichment failed: {exc}", exc_info=True)
                    return ("census_demographics", None)

            def _call_usajobs() -> "tuple[str, Optional[dict]]":
                """Fetch USAJobs federal listings (only for government/public sector)."""
                if not (
                    _api_integrations_available and _api_usajobs and _roles_normalized
                ):
                    return ("usajobs_federal", None)
                # Only call for government-adjacent industries
                _industry_lower = str(data.get("industry") or "").lower()
                _gov_keywords = {
                    "government",
                    "public",
                    "federal",
                    "defense",
                    "military",
                    "healthcare",
                    "education",
                }
                if not any(kw in _industry_lower for kw in _gov_keywords):
                    return ("usajobs_federal", None)
                per_role: list[dict[str, Any]] = []
                for _ri in _roles_normalized[:3]:
                    try:
                        count = _api_usajobs.get_job_count(_ri["title"])
                        if count:
                            per_role.append(
                                {"role": _ri["title"], "federal_openings": count}
                            )
                    except (
                        urllib.error.URLError,
                        OSError,
                        ValueError,
                        TypeError,
                    ) as exc:
                        logger.error(
                            f"USAJobs enrichment for role '{_ri['title']}' failed: {exc}",
                            exc_info=True,
                        )
                if not per_role:
                    return ("usajobs_federal", None)
                total = sum(r.get("federal_openings") or 0 for r in per_role)
                logger.info(
                    f"Enriched /api/generate with USAJobs data for {len(per_role)} role(s) ({total} openings)"
                )
                return ("usajobs_federal", {"total": total, "per_role": per_role})

            def _call_knowledge_graph() -> "tuple[str, Optional[dict]]":
                """Fetch Google Knowledge Graph entity data for company and roles."""
                if not _has_knowledge_graph:
                    return ("knowledge_graph", None)
                try:
                    result: dict[str, Any] = {}
                    _client_name = data.get("client_name") or ""
                    if _client_name:
                        company_kg = enrich_company(_client_name)
                        if (
                            company_kg
                            and isinstance(company_kg, dict)
                            and company_kg.get("name")
                        ):
                            result["company"] = company_kg
                    # Enrich first role for occupation context
                    if _roles_normalized:
                        role_kg = enrich_job_title(_roles_normalized[0]["title"])
                        if (
                            role_kg
                            and isinstance(role_kg, dict)
                            and role_kg.get("name")
                        ):
                            result["primary_role"] = role_kg
                    # Enrich first location for geo context
                    if isinstance(_locs_for_api, list) and _locs_for_api:
                        _loc_val = _locs_for_api[0]
                        _loc_str_kg = (
                            _loc_val if isinstance(_loc_val, str) else str(_loc_val)
                        )
                        loc_kg = enrich_location(_loc_str_kg)
                        if loc_kg and isinstance(loc_kg, dict) and loc_kg.get("name"):
                            result["primary_location"] = loc_kg
                    if not result:
                        return ("knowledge_graph", None)
                    logger.info(
                        f"Enriched /api/generate with Knowledge Graph data ({list(result.keys())})"
                    )
                    return ("knowledge_graph", result)
                except (urllib.error.URLError, OSError, ValueError, TypeError) as exc:
                    logger.error(
                        f"Knowledge Graph enrichment failed: {exc}", exc_info=True
                    )
                    return ("knowledge_graph", None)

            def _call_google_maps() -> "tuple[str, Optional[dict]]":
                """Geocode campaign locations for coordinate/timezone enrichment."""
                if not (_has_google_maps and _locs_for_api):
                    return ("google_maps_geo", None)
                try:
                    locs_to_geocode: list[str] = []
                    for _loc in (
                        _locs_for_api if isinstance(_locs_for_api, list) else []
                    ):
                        if isinstance(_loc, str) and _loc.strip():
                            locs_to_geocode.append(_loc.strip())
                        elif isinstance(_loc, dict):
                            _parts_g = [
                                _loc.get("city") or "",
                                _loc.get("state") or "",
                                _loc.get("country") or "",
                            ]
                            _joined = ", ".join(filter(None, _parts_g))
                            if _joined:
                                locs_to_geocode.append(_joined)
                    if not locs_to_geocode:
                        return ("google_maps_geo", None)
                    geocoded = batch_geocode(locs_to_geocode[:5])  # Cap at 5 locations
                    if geocoded and isinstance(geocoded, list):
                        valid = [
                            g for g in geocoded if isinstance(g, dict) and g.get("lat")
                        ]
                        if valid:
                            logger.info(
                                f"Enriched /api/generate with Google Maps geocoding for {len(valid)} location(s)"
                            )
                            return ("google_maps_geo", {"locations": valid})
                    return ("google_maps_geo", None)
                except (urllib.error.URLError, OSError, ValueError, TypeError) as exc:
                    logger.error(f"Google Maps enrichment failed: {exc}", exc_info=True)
                    return ("google_maps_geo", None)

            # -- Dispatch all enrichment tasks concurrently (20s combined timeout) --
            _enrich_tasks = [
                _call_enrich_data,
                _call_adzuna,
                _call_bls,
                _call_fred,
                _call_jobspy,
                _call_onet,
                _call_jooble,
                _call_bea,
                _call_census,
                _call_usajobs,
                _call_knowledge_graph,
                _call_google_maps,
            ]
            # ── TIMEOUT CHECK: Before enrichment phase (expensive) ──
            if _gen_timeout_flag.is_set():
                _gen_timer.cancel()
                elapsed = time.time() - _gen_start
                logger.error(
                    f"Generation timeout (before enrichment phase) after {elapsed:.1f}s"
                )
                self._send_error(
                    f"Generation timed out after {_GENERATE_TIMEOUT_SECONDS} seconds. Please simplify your request.",
                    "TIMEOUT",
                    504,
                )
                return

            _parallel_t0 = time.time()

            # CRITICAL: Do NOT use `with ThreadPoolExecutor() as pool:` context manager.
            # The context manager calls pool.shutdown(wait=True) on exit, which blocks
            # until ALL threads complete -- even after TimeoutError is raised.
            # Threads stuck on semaphore.acquire() or slow API calls will block the
            # entire response. Use manual pool with shutdown(wait=False) instead.
            _enrich_pool = ThreadPoolExecutor(
                max_workers=8, thread_name_prefix="gen-enrich"
            )
            try:
                with _span_fn(
                    "enrich_parallel",
                    "Parallel API enrichment (12 sources, 20s timeout)",
                ):
                    futures = {
                        _enrich_pool.submit(fn): fn.__doc__ or fn.__name__
                        for fn in _enrich_tasks
                    }
                    for future in as_completed(futures, timeout=20):
                        try:
                            name, result = future.result()
                            if name == "api_enrichment":
                                if result:
                                    enriched = result
                            else:
                                if result is not None:
                                    _market_ctx[name] = result
                        except Exception as exc:
                            logger.warning("Generate enrichment task error: %s", exc)
            except TimeoutError:
                elapsed_time: float = time.time() - _parallel_t0
                logger.warning(
                    "Parallel enrichment hit 20s timeout after %.1fs -- "
                    "continuing with partial data (not aborting)",
                    elapsed_time,
                )
                # Cancel pending futures (running threads can't be cancelled)
                for future in futures:
                    future.cancel()
                # DON'T abort -- continue with whatever data we have
            except Exception as exc:
                logger.error("Parallel enrichment pool failed: %s", exc, exc_info=True)
            finally:
                # shutdown(wait=False) exits immediately without blocking
                _enrich_pool.shutdown(wait=False, cancel_futures=True)

            # Enrichment complete (or timed out with partial data) -- continue

            data["_enriched"] = enriched
            if _market_ctx:
                data["market_context"] = _market_ctx

            # ── Fix 25 + S49 Issue 17: Detect enrichment quality degradation ──
            # Scale warning text with confidence level so PPT disclaimer
            # is stronger when data quality is worse.
            _enr_sync = data.get("_enriched") or {}
            _enr_sync_summary = (
                _enr_sync.get("enrichment_summary", {})
                if isinstance(_enr_sync, dict)
                else {}
            )
            _enr_sync_succeeded = _enr_sync_summary.get("apis_succeeded") or []
            _enr_sync_confidence = float(
                _enr_sync_summary.get("confidence_score", 0)
                if isinstance(_enr_sync_summary, dict)
                else 0
            )
            if not _enr_sync or (
                isinstance(_enr_sync_succeeded, list) and len(_enr_sync_succeeded) < 3
            ):
                # S49: Tiered warning based on confidence score
                if _enr_sync_confidence < 0.40:
                    data["quality_warning"] = (
                        "Minimal data available. Plan is heavily estimated. "
                        "Recommend API enrichment before executing."
                    )
                elif _enr_sync_confidence < 0.60:
                    data["quality_warning"] = (
                        "Limited data available. This plan uses estimated benchmarks. "
                        "Results may vary significantly."
                    )
                else:
                    data["quality_warning"] = (
                        "Some market data sources were unavailable. "
                        "Plan uses benchmark estimates."
                    )
                logger.warning(
                    "Enrichment quality degraded: %d APIs succeeded (min 3), confidence=%.2f",
                    (
                        len(_enr_sync_succeeded)
                        if isinstance(_enr_sync_succeeded, list)
                        else 0
                    ),
                    _enr_sync_confidence,
                )

            logger.info(
                "Parallel enrichment completed in %.2fs (sources: %s)",
                time.time() - _parallel_t0,
                ", ".join(
                    ["api_enrichment"] * bool(enriched) + list(_market_ctx.keys())
                )
                or "none",
            )

            # ── PostHog: Track enrichment pipeline completion ──
            _enrich_apis_used: list[str] = []
            if data.get("_enriched"):
                _enrich_apis_used.append("api_enrichment")
            if data.get("market_context"):
                _enrich_apis_used.extend(
                    list((data.get("market_context") or {}).keys())
                )
            if _enrich_apis_used:
                self._ph_track(
                    "api_enrichment_complete",
                    {
                        "apis_called": _enrich_apis_used,
                        "total_latency_ms": round((time.time() - _gen_start) * 1000),
                    },
                )

            # ── TIMEOUT CHECK: Before knowledge base loading ──
            if _gen_timeout_flag.is_set():
                _gen_timer.cancel()
                elapsed = time.time() - _gen_start
                logger.error(
                    f"Generation timeout (after enrichment) after {elapsed:.1f}s"
                )
                self._send_error(
                    f"Generation timed out after {_GENERATE_TIMEOUT_SECONDS} seconds. Please simplify your request.",
                    "TIMEOUT",
                    504,
                )
                return

            # ── Phase 2: Load Knowledge Base ──
            with _span_fn("kb.load", "Load knowledge base"):
                kb = load_knowledge_base()
            # Enrich KB with vendor/publisher profiles for channel strategy
            if _supabase_data_available:
                try:
                    _sb_vendors = get_vendor_profiles()
                    if _sb_vendors:
                        kb["_supabase_vendor_profiles"] = _sb_vendors
                        data["_vendor_profiles"] = _sb_vendors
                except Exception as _vp_err:
                    logger.warning(
                        "Vendor profiles enrichment failed (non-fatal): %s",
                        _vp_err,
                    )
            data["_knowledge_base"] = kb  # Pass KB to PPT for fallback data

            # ── TIMEOUT CHECK: Before data synthesis ──
            if _gen_timeout_flag.is_set():
                _gen_timer.cancel()
                elapsed = time.time() - _gen_start
                logger.error(
                    f"Generation timeout (before synthesis) after {elapsed:.1f}s"
                )
                self._send_error(
                    f"Generation timed out after {_GENERATE_TIMEOUT_SECONDS} seconds. Please simplify your request.",
                    "TIMEOUT",
                    504,
                )
                return

            # ── Phase 3: Data Synthesis ──
            _t_synth_sync = time.time()
            if data_synthesize is not None:
                try:
                    with _span_fn("synthesize", "Data synthesis (fuse enriched + KB)"):
                        synthesized = data_synthesize(enriched, kb, data)
                    data["_synthesized"] = synthesized
                    logger.info(
                        "Data synthesis complete: %s",
                        (
                            list(synthesized.keys())
                            if isinstance(synthesized, dict)
                            else "N/A"
                        ),
                    )
                except Exception as e:
                    logger.warning("Data synthesis failed (non-fatal): %s", e)
                    data["_synthesized"] = {}
            else:
                data["_synthesized"] = {}

            logger.info(
                "PERF [sync]: Synthesis took %.2fs", time.time() - _t_synth_sync
            )

            # ── Copy geopolitical context from enriched to synthesized ──
            # PPT/Excel read from _synthesized, but geopolitical_context is produced
            # by api_enrichment.enrich_data() and stored in _enriched. Bridge the gap.
            if isinstance(enriched, dict) and isinstance(
                data.get("_synthesized"), dict
            ):
                _geo_from_enriched = enriched.get("geopolitical_context")
                if _geo_from_enriched and isinstance(_geo_from_enriched, dict):
                    data["_synthesized"]["geopolitical_context"] = _geo_from_enriched

            # ── NAICS-based Industry Classification ──
            # Classify industry using the NAICS engine, supporting both legacy keys
            # (from frontend dropdown) and free-text industry names
            industry_raw = data.get("industry") or ""
            company_name = (
                data.get("client_name") or "" or data.get("company_name") or ""
            )
            roles_list = data.get("target_roles") or data.get("roles") or []
            if isinstance(roles_list, str):
                roles_list = [r.strip() for r in roles_list.split(",") if r.strip()]
            industry_profile = classify_industry(industry_raw, company_name, roles_list)

            # Set the industry key to the legacy key for backward compatibility with
            # research.py, channels_db.json, and all internal lookup tables
            data["industry"] = industry_profile.get("legacy_key", "general_entry_level")

            # Set the display label from NAICS classification (prefer frontend-provided label)
            if not data.get("industry_label"):
                data["industry_label"] = industry_profile["sector"]

            # Store talent profile and BLS sector for downstream use
            data["talent_profile"] = industry_profile["talent_profile"]
            data["bls_sector"] = industry_profile["bls_sector"]
            data["naics_code"] = industry_profile.get("naics", "00")

            # ── Phase 4: Budget Allocation ──
            _t_budget_sync = time.time()
            if calculate_budget_allocation is not None:
                try:
                    _ind_key_ba = data.get("industry", "general_entry_level")
                    _DEFAULT_ALLOC_BA = {
                        "programmatic_dsp": 35,
                        "global_boards": 20,
                        "niche_boards": 15,
                        "social_media": 12,
                        "regional_boards": 8,
                        "employer_branding": 5,
                        "apac_regional": 3,
                        "emea_regional": 2,
                    }
                    # Use the expanded 17-industry profiles from ppt_generator if available,
                    # falling back to the default allocation.
                    # CRITICAL: dict.copy() to avoid mutating the shared INDUSTRY_ALLOC_PROFILES
                    # (the region-stripping .pop() calls below would permanently corrupt the shared dict)
                    if INDUSTRY_ALLOC_PROFILES is not None:
                        channel_pcts = dict(
                            INDUSTRY_ALLOC_PROFILES.get(_ind_key_ba, _DEFAULT_ALLOC_BA)
                        )
                    else:
                        # Inline fallback if ppt_generator is not importable (6 industries)
                        _INDUSTRY_ALLOC_BA_FALLBACK = {
                            "healthcare_medical": {
                                "programmatic_dsp": 22,
                                "global_boards": 15,
                                "niche_boards": 30,
                                "social_media": 10,
                                "regional_boards": 10,
                                "employer_branding": 8,
                                "apac_regional": 3,
                                "emea_regional": 2,
                            },
                            "tech_engineering": {
                                "programmatic_dsp": 30,
                                "global_boards": 15,
                                "niche_boards": 20,
                                "social_media": 18,
                                "regional_boards": 5,
                                "employer_branding": 7,
                                "apac_regional": 3,
                                "emea_regional": 2,
                            },
                            "finance_banking": {
                                "programmatic_dsp": 25,
                                "global_boards": 18,
                                "niche_boards": 25,
                                "social_media": 10,
                                "regional_boards": 7,
                                "employer_branding": 10,
                                "apac_regional": 3,
                                "emea_regional": 2,
                            },
                            "retail_consumer": {
                                "programmatic_dsp": 38,
                                "global_boards": 22,
                                "niche_boards": 8,
                                "social_media": 20,
                                "regional_boards": 7,
                                "employer_branding": 3,
                                "apac_regional": 1,
                                "emea_regional": 1,
                            },
                            "general_entry_level": {
                                "programmatic_dsp": 40,
                                "global_boards": 22,
                                "niche_boards": 8,
                                "social_media": 15,
                                "regional_boards": 10,
                                "employer_branding": 3,
                                "apac_regional": 1,
                                "emea_regional": 1,
                            },
                            # S49 Issue 15: Transportation/logistics/trucking
                            # niche CDL boards (CDLjobs, TruckersReport, DriveMyWay) dominate
                            "logistics_supply_chain": {
                                "programmatic_dsp": 25,
                                "global_boards": 15,
                                "niche_boards": 35,
                                "social_media": 10,
                                "regional_boards": 10,
                                "employer_branding": 5,
                                "apac_regional": 0,
                                "emea_regional": 0,
                            },
                        }
                        channel_pcts = dict(
                            _INDUSTRY_ALLOC_BA_FALLBACK.get(
                                _ind_key_ba, _DEFAULT_ALLOC_BA
                            )
                        )

                    # Strip APAC/EMEA channels based on target_region or location inference
                    _target_region = data.get("target_region") or "us_only"
                    _locs_raw = data.get("locations") or []
                    _all_us = _target_region == "us_only"
                    if not _all_us and _target_region not in (
                        "global",
                        "emea",
                        "apac",
                        "custom",
                    ):
                        # Infer from locations if region not explicitly set
                        _all_us = True
                        if _locs_raw:
                            for _loc in (
                                _locs_raw
                                if isinstance(_locs_raw, list)
                                else [_locs_raw]
                            ):
                                _loc_str = str(
                                    _loc.get("country")
                                    if isinstance(_loc, dict)
                                    else _loc
                                ).lower()
                                if _loc_str and _loc_str not in (
                                    "us",
                                    "usa",
                                    "united states",
                                    "",
                                ):
                                    _all_us = False
                                    break

                    if isinstance(channel_pcts, dict):
                        if _all_us or _target_region == "us_only":
                            # US-only: remove APAC/EMEA channels
                            _intl_pct = channel_pcts.pop(
                                "apac_regional", 0
                            ) + channel_pcts.pop("emea_regional", 0)
                            if _intl_pct > 0:
                                _top_ch = max(
                                    channel_pcts, key=lambda k: channel_pcts[k]
                                )
                                channel_pcts[_top_ch] = (
                                    channel_pcts[_top_ch] + _intl_pct
                                )
                                logger.info(
                                    f"US-only plan (region={_target_region}): redistributed {_intl_pct}%% from APAC/EMEA to {_top_ch}"
                                )
                        elif _target_region == "emea":
                            # EMEA-only: remove APAC channels, boost EMEA
                            _apac_pct = channel_pcts.pop("apac_regional", 0)
                            if _apac_pct > 0:
                                channel_pcts["emea_regional"] = (
                                    channel_pcts.get("emea_regional", 0) + _apac_pct
                                )
                                logger.info(
                                    f"EMEA plan: redistributed {_apac_pct}%% from APAC to EMEA"
                                )
                        elif _target_region == "apac":
                            # APAC-only: remove EMEA channels, boost APAC
                            _emea_pct = channel_pcts.pop("emea_regional", 0)
                            if _emea_pct > 0:
                                channel_pcts["apac_regional"] = (
                                    channel_pcts.get("apac_regional", 0) + _emea_pct
                                )
                                logger.info(
                                    f"APAC plan: redistributed {_emea_pct}%% from EMEA to APAC"
                                )

                    # Parse budget to float — uses shared_utils.parse_budget (single source of truth)
                    _bstr_ba = str(
                        data.get("budget") or "" or data.get("budget_range") or "" or ""
                    ).strip()
                    if not _bstr_ba:
                        _bstr_ba = str(data.get("budget_range") or "" or "").strip()
                    _bval_ba = parse_budget(_bstr_ba)
                    logger.info(
                        "Budget allocation: parsed '%s' -> $%s",
                        _bstr_ba,
                        f"{_bval_ba:,.2f}",
                    )

                    # Build role dicts from string list
                    _roles_raw = data.get("target_roles") or data.get("roles") or []
                    _roles_for_ba = []
                    for r in (_roles_raw if isinstance(_roles_raw, list) else []):
                        if isinstance(r, str):
                            _roles_for_ba.append(
                                {
                                    "title": r,
                                    "count": 1,
                                    "tier": data.get("_role_tiers", {})
                                    .get(r, {})
                                    .get("tier", "Professional"),
                                }
                            )
                        elif isinstance(r, dict):
                            _roles_for_ba.append(
                                {
                                    "title": r.get("title") or "",
                                    "count": int(r.get("count", 1)),
                                    "tier": r.get("_tier", "Professional"),
                                }
                            )

                    # Build location dicts from string list
                    _locs_raw = data.get("locations") or []
                    _locs_for_ba = []
                    for loc in (_locs_raw if isinstance(_locs_raw, list) else []):
                        if isinstance(loc, str):
                            parts = [p.strip() for p in loc.split(",")]
                            _locs_for_ba.append(
                                {
                                    "city": parts[0] if parts else loc,
                                    "state": parts[1] if len(parts) > 1 else "",
                                    "country": (
                                        parts[2]
                                        if len(parts) > 2
                                        else parts[-1] if len(parts) > 1 else "US"
                                    ),
                                }
                            )
                        elif isinstance(loc, dict):
                            _locs_for_ba.append(
                                {
                                    "city": loc.get("city") or "",
                                    "state": loc.get("state") or "",
                                    "country": loc.get("country") or "",
                                }
                            )

                    synthesized_for_ba = data.get("_synthesized", {})
                    enriched_for_ba = data.get("_enriched", {})
                    # Merge raw enrichment + synthesized so budget engine can
                    # access both raw API keys (google_ads_data, teleport_data)
                    # and synthesized fusion keys (ad_platform_analysis, etc.)
                    merged_for_ba = {}
                    if isinstance(enriched_for_ba, dict):
                        merged_for_ba.update(enriched_for_ba)
                    if isinstance(synthesized_for_ba, dict):
                        merged_for_ba.update(synthesized_for_ba)
                    budget_result = calculate_budget_allocation(
                        total_budget=_bval_ba,
                        roles=_roles_for_ba,
                        locations=_locs_for_ba,
                        industry=data.get("industry", "General"),
                        channel_percentages=channel_pcts,
                        synthesized_data=merged_for_ba,
                        knowledge_base=kb,
                        collar_type=data.get("_collar_type") or "",
                        campaign_start_month=int(
                            data.get("campaign_start_month") or 0 or 0
                        ),
                    )
                    data["_budget_allocation"] = budget_result
                    logger.info(
                        "Budget allocation complete: %s",
                        (
                            list(budget_result.keys())
                            if isinstance(budget_result, dict)
                            else "N/A"
                        ),
                    )
                except Exception as e:
                    logger.warning("Budget allocation failed (non-fatal): %s", e)
                    data["_budget_allocation"] = {}
            else:
                data["_budget_allocation"] = {}

            logger.info(
                "PERF [sync]: Budget allocation took %.2fs",
                time.time() - _t_budget_sync,
            )

            # ── Inject SlotOps LinkedIn Benchmarks (sync path) ──
            # Pulls from 108K-job dataset: country apply rates, EA vs ATS, best days
            try:
                from slotops_engine import get_linkedin_benchmarks_for_plan

                _slotops_country = "United States"
                _locs_for_slotops = data.get("locations") or []
                if _locs_for_slotops:
                    _first_loc = (
                        _locs_for_slotops[0]
                        if isinstance(_locs_for_slotops, list)
                        else _locs_for_slotops
                    )
                    if isinstance(_first_loc, dict):
                        _slotops_country = _first_loc.get("country") or "United States"
                    elif isinstance(_first_loc, str):
                        parts = [p.strip() for p in _first_loc.split(",")]
                        if len(parts) >= 2:
                            _slotops_country = parts[-1]
                _li_benchmarks = get_linkedin_benchmarks_for_plan(_slotops_country)
                if _li_benchmarks:
                    data["_slotops_linkedin_benchmarks"] = _li_benchmarks
                    logger.info(
                        "SlotOps LinkedIn benchmarks injected for %s (sample=%s)",
                        _slotops_country,
                        _li_benchmarks.get("sample_size", 0),
                    )
            except ImportError:
                logger.debug(
                    "slotops_engine not available -- skipping LinkedIn benchmarks"
                )
            except Exception as _slotops_err:
                logger.warning(
                    "SlotOps LinkedIn benchmarks failed (non-fatal): %s", _slotops_err
                )

            # ── Creative Quality Score (P1-16) ──
            try:
                from creative_quality_score import score_creative_quality

                data["_creative_quality_score"] = score_creative_quality(data)
                logger.info(
                    "Creative quality score: %s (%s)",
                    data["_creative_quality_score"].get("score"),
                    data["_creative_quality_score"].get("grade"),
                )
            except ImportError:
                logger.debug("creative_quality_score not available -- skipping")
            except Exception as _cqs_err:
                logger.warning(
                    "Creative quality score failed (non-fatal): %s", _cqs_err
                )

            # ── Gold Standard Quality Gates ──
            # Apply all 7 quality gates: city-level data, clearance segmentation,
            # competitor mapping, difficulty framework, channel strategy with splits,
            # multi-tier budget breakdowns, and activation event calendars.
            # Wrapped in ThreadPoolExecutor with 30s timeout to prevent hangs.
            _t_gs_sync = time.time()
            try:
                from gold_standard import apply_all_quality_gates

                _gs_pool = ThreadPoolExecutor(max_workers=1)
                _gs_future = _gs_pool.submit(apply_all_quality_gates, data)
                try:
                    _gs_result = _gs_future.result(timeout=30)
                    logger.info(
                        "Gold Standard gates complete: %d gates produced data",
                        len(_gs_result),
                    )
                except TimeoutError:
                    logger.warning(
                        "Gold Standard gates timed out after 30s -- continuing with partial data"
                    )
                    _gs_future.cancel()
                    data["_gold_standard"] = data.get("_gold_standard") or {}
                finally:
                    _gs_pool.shutdown(wait=False)
            except ImportError:
                logger.warning(
                    "gold_standard module not available -- skipping quality gates"
                )
                data["_gold_standard"] = {}
            except Exception as _gs_err:
                logger.error(
                    "Gold Standard gates failed (non-fatal): %s", _gs_err, exc_info=True
                )
                data["_gold_standard"] = {}
            logger.info(
                "PERF [sync]: Gold Standard gates took %.2fs", time.time() - _t_gs_sync
            )

            # ── Gemini/LLM verification of plan data (same as async path) ──
            # Wrapped in ThreadPoolExecutor with 30s timeout to prevent hangs.
            try:
                _vf_pool = ThreadPoolExecutor(max_workers=1)
                _vf_future = _vf_pool.submit(_verify_plan_data, data)
                try:
                    data["_verification"] = _vf_future.result(timeout=30)
                except TimeoutError:
                    logger.warning(
                        "Plan data verification timed out after 30s -- skipping"
                    )
                    _vf_future.cancel()
                    data["_verification"] = {
                        "status": "skipped",
                        "reason": "verification_timeout",
                    }
                finally:
                    _vf_pool.shutdown(wait=False)
            except Exception:
                data["_verification"] = {
                    "status": "skipped",
                    "reason": "verification_error",
                }

            # PostHog: track plan generation started (funnel entry)
            try:
                from posthog_tracker import track_plan_started

                track_plan_started(
                    email=data.get("requester_email", "anonymous"),
                    client=data.get("client_name") or "",
                    industry=data.get("industry") or "",
                    roles_count=len(
                        data.get("target_roles") or data.get("roles") or []
                    ),
                )
            except Exception:
                pass

            # ── Inject international benchmarks for non-US plans (sync path) ──
            _target_region_sync = data.get("target_region") or "us_only"
            if _target_region_sync in ("global", "emea", "apac", "custom"):
                try:
                    _kb_intl_sync = load_knowledge_base().get(
                        "international_benchmarks", {}
                    )
                    if _kb_intl_sync:
                        _intl_countries_s = _kb_intl_sync.get("countries", {})
                        _intl_regions_s = _kb_intl_sync.get("regions", {})
                        _relevant_countries_s: dict = {}
                        if _target_region_sync == "emea":
                            _region_keys_s = ["emea"]
                        elif _target_region_sync == "apac":
                            _region_keys_s = ["apac"]
                        else:
                            _region_keys_s = list(_intl_regions_s.keys())
                        for _rk_s in _region_keys_s:
                            _rc_s = _intl_regions_s.get(_rk_s, {}).get("countries", [])
                            for _ck_s in _rc_s:
                                if _ck_s in _intl_countries_s:
                                    _relevant_countries_s[_ck_s] = _intl_countries_s[
                                        _ck_s
                                    ]
                        data["_intl_benchmarks"] = {
                            "countries": _relevant_countries_s,
                            "regions": {
                                k: v
                                for k, v in _intl_regions_s.items()
                                if k in _region_keys_s
                            },
                            "source": "international_benchmarks_2026 (38 countries)",
                        }
                        logger.info(
                            "Injected intl benchmarks (sync): %d countries for region=%s",
                            len(_relevant_countries_s),
                            _target_region_sync,
                        )
                except Exception as _intl_err_s:
                    logger.error(
                        "Failed to inject intl benchmarks (sync): %s",
                        _intl_err_s,
                        exc_info=True,
                    )

            # ── Cross-validation post-processing (sync path) ──
            # Detects internal inconsistencies in the enriched plan and
            # auto-corrects where safe.  Non-fatal: failures here never
            # block Excel/PPT generation.
            try:
                from plan_validator import validate_plan

                _val_result_sync = validate_plan(data)
                logger.info(
                    "Plan validation (sync): %d finding(s), %d auto-corrected",
                    _val_result_sync.get("total_findings", 0),
                    _val_result_sync.get("auto_corrections", 0),
                )
            except ImportError:
                logger.debug("plan_validator not available -- skipping")
            except Exception as _val_err_sync:
                logger.warning(
                    "Plan validation failed (sync, non-fatal): %s",
                    _val_err_sync,
                    exc_info=True,
                )
                data["_validation"] = {
                    "findings": [],
                    "checks_run": 0,
                    "checks_failed": 1,
                    "error": str(_val_err_sync),
                }

            # ── Excel generation phase (timeout timer already started at request entry) ──
            try:
                # Check timeout before heavy work
                if _gen_timeout_flag.is_set():
                    raise _GenerationTimeoutError(
                        "Generation exceeded 4-minute timeout"
                    )
                with _span_fn("generate.excel", "Generate Excel media plan"):
                    # Use v2 consolidated 4-sheet generator with legacy fallback
                    if generate_excel_v2 is not None:
                        try:
                            excel_bytes = generate_excel_v2(
                                data,
                                research_mod=research,
                                load_kb_fn=load_knowledge_base,
                                classify_tier_fn=classify_role_tier,
                                fetch_logo_fn=fetch_client_logo,
                            )
                            logger.info(
                                "Excel generated via excel_v2 (4-sheet consolidated)"
                            )
                        except Exception as v2_err:
                            logger.warning(
                                "excel_v2 failed, falling back to legacy: %s", v2_err
                            )
                            if generate_excel is not None:
                                excel_bytes = generate_excel(data)
                            else:
                                raise RuntimeError(
                                    "Both excel_v2 and legacy generator unavailable"
                                ) from v2_err
                    else:
                        if generate_excel is not None:
                            excel_bytes = generate_excel(data)
                        else:
                            raise RuntimeError(
                                "No Excel generator available (excel_v2 and legacy both missing)"
                            )
            except _GenerationTimeoutError:
                _gen_timer.cancel()
                elapsed = time.time() - _gen_start
                logger.error(
                    f"Generation timeout after {elapsed:.1f}s for client={data.get('client_name')}"
                )
                self._send_error(
                    f"Generation timed out after {_GENERATE_TIMEOUT_SECONDS} seconds. Please simplify your request.",
                    "TIMEOUT",
                    504,
                )
                return
            except Exception as e:
                _gen_timer.cancel()
                import traceback

                # Check if we hit the timeout during generation
                if _gen_timeout_flag.is_set():
                    elapsed = time.time() - _gen_start
                    logger.error(
                        f"Generation timeout (detected in except) after {elapsed:.1f}s"
                    )
                    self._send_error(
                        f"Generation timed out after {_GENERATE_TIMEOUT_SECONDS} seconds. Please simplify your request.",
                        "TIMEOUT",
                        504,
                    )
                    return

                tb = traceback.format_exc()
                logger.error("Excel generation error: %s", tb)
                generation_time = time.time() - _gen_start
                log_request(
                    data, "error", generation_time=generation_time, error_msg=str(e)
                )
                try:
                    from email_alerts import send_generation_failure_alert

                    send_generation_failure_alert(
                        client_name=data.get("client_name", "Unknown"),
                        error=str(e),
                        request_id=data.get("_request_id") or "",
                    )
                except Exception:
                    pass
                # PostHog: track failed generation (both legacy + funnel event)
                try:
                    from posthog_tracker import track_plan_failed, capture

                    track_plan_failed(
                        email=data.get("requester_email", "anonymous"),
                        client=data.get("client_name") or "",
                        error=str(e),
                    )
                    capture(
                        "media_plan_failed",
                        distinct_id=data.get("requester_email", "anonymous"),
                        properties={
                            "client_name": data.get("client_name") or "",
                            "error": str(e)[:200],
                            "source": "server",
                        },
                    )
                except Exception:
                    pass
                self._send_error(
                    "Generation failed. Please check your inputs and try again.",
                    "SERVER_ERROR",
                    500,
                )
                return
            finally:
                _gen_timer.cancel()

            # Sanitize client_name to ASCII-safe characters (prevents CJK/Unicode crashes in filenames/headers)
            _raw_client = data.get("client_name") or "Client"
            client_name = re.sub(
                r"_+",
                "_",  # S27: Collapse multiple underscores
                re.sub(r"[^a-zA-Z0-9_\-]", "_", _raw_client),
            ).strip("_")
            # S47: Build descriptive filename with industry, location, and date
            _fn_industry = re.sub(
                r"[^a-zA-Z0-9]", "_", (data.get("industry") or "").strip().title()
            ).strip("_")[:30]
            _fn_locations = data.get("locations") or []
            _fn_location = ""
            if isinstance(_fn_locations, list) and _fn_locations:
                _first_loc = (
                    _fn_locations[0]
                    if isinstance(_fn_locations[0], str)
                    else str(_fn_locations[0])
                )
                _fn_location = re.sub(r"[^a-zA-Z0-9]", "_", _first_loc.strip()).strip(
                    "_"
                )[:20]
            elif isinstance(_fn_locations, str) and _fn_locations.strip():
                _fn_location = re.sub(
                    r"[^a-zA-Z0-9]", "_", _fn_locations.strip()
                ).strip("_")[:20]
            _fn_date = datetime.datetime.now().strftime("%Y-%m-%d")
            # Build: "Media_Plan_Healthcare_US_2026-04-07" or "ClientName_Media_Plan_2026-04-07"
            _fn_parts = ["Media_Plan"]
            if _fn_industry:
                _fn_parts.append(_fn_industry)
            if _fn_location:
                _fn_parts.append(_fn_location)
            _fn_parts.append(_fn_date)
            _descriptive_filename = "_".join(_fn_parts)
            # Collapse multiple underscores
            _descriptive_filename = re.sub(r"_+", "_", _descriptive_filename)

            # Store plan results JSON for on-screen dashboard
            _plan_id = uuid.uuid4().hex[:12]
            _store_plan_result(_plan_id, data)

            # ── Async Google Sheets export (non-blocking) ──
            # Creates a cloud-hosted Google Sheet copy of the plan in the
            # background.  The URL is stored in _plan_results_store for
            # later retrieval via /api/plan/sheets-url?plan_id=xxx.
            def _create_sheets_async(
                plan_data_copy: dict,
                plan_id_ref: str,
            ) -> None:
                try:
                    from sheets_export import export_media_plan

                    sheet_url = export_media_plan(plan_data_copy)
                    if sheet_url:
                        with _plan_results_lock:
                            entry = _plan_results_store.get(plan_id_ref)
                            if entry:
                                entry["sheets_url"] = sheet_url
                        logger.info(
                            "Async Google Sheets created for plan %s: %s",
                            plan_id_ref,
                            sheet_url,
                        )
                except Exception as exc:
                    logger.error(
                        "Async Sheets export failed for plan %s: %s",
                        plan_id_ref,
                        exc,
                        exc_info=True,
                    )

            try:
                submit_background_task(
                    _create_sheets_async,
                    dict(data),  # shallow copy to avoid mutation
                    _plan_id,
                    task_name=f"sheets-export-{_plan_id}",
                    timeout=45,
                )
            except Exception as _sheets_bg_err:
                logger.warning(
                    "Async Sheets export submission failed (non-fatal): %s",
                    _sheets_bg_err,
                )

            # ── Store plan in BigQuery for historical analytics (non-blocking) ──
            if _has_bigquery:
                try:
                    bq_store_plan(dict(data))
                except Exception as _bq_err:
                    logger.warning("BigQuery plan storage failed: %s", _bq_err)

            # Generate Strategy PPT deck (7-tier fallback via DeckGenerator)
            pptx_bytes = None
            pptx_warning = None
            pptx_provider = None
            # Check for ?deck_tier=premium query param to force Gamma (tier 2)
            _deck_tier_param = (
                _qs.get("deck_tier", [None])[0] if "_qs" in dir() else None
            )
            if _deck_tier_param is None:
                _parsed_qs = urllib.parse.parse_qs(
                    urllib.parse.urlparse(self.path).query
                )
                _deck_tier_param = _parsed_qs.get("deck_tier", [None])[0]
            _force_tier = None
            if _deck_tier_param == "premium":
                _force_tier = "gamma"
            elif _deck_tier_param and _deck_tier_param != "auto":
                _force_tier = _deck_tier_param

            if _deck_generator is not None:
                try:
                    with _span_fn(
                        "generate.pptx", "Generate strategy PPT deck (multi-tier)"
                    ):
                        pptx_bytes, pptx_provider = _deck_generator.generate(
                            data, force_tier=_force_tier
                        )
                    logger.info(
                        "PPT generated via %s: %d bytes",
                        pptx_provider,
                        len(pptx_bytes),
                    )
                except Exception as e:
                    logger.error("DeckGenerator error: %s", e, exc_info=True)
                    pptx_bytes = None
                    pptx_warning = "Strategy deck (PPT) could not be generated. Excel plan is included."
            elif generate_pptx is not None:
                # Direct fallback if deck_generator module not loaded
                try:
                    with _span_fn("generate.pptx", "Generate strategy PPT deck"):
                        pptx_bytes = generate_pptx(data)
                    pptx_provider = "pptx"
                    logger.info("PPT generated (direct): %d bytes", len(pptx_bytes))
                except Exception as e:
                    logger.error("PPT generation error: %s", e, exc_info=True)
                    pptx_bytes = None
                    pptx_warning = "Strategy deck (PPT) could not be generated. Excel plan is included."
            else:
                logger.info(
                    "PPT generation skipped: neither deck_generator nor ppt_generator available"
                )

            if pptx_bytes:
                # Bundle both files in a ZIP
                zip_buffer = io.BytesIO()
                with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zf:
                    zf.writestr(f"{_descriptive_filename}.xlsx", excel_bytes)
                    zf.writestr(f"{_descriptive_filename}_Deck.pptx", pptx_bytes)
                zip_bytes = zip_buffer.getvalue()

                # Save a copy for the document repository
                doc_filename = None
                try:
                    docs_dir = os.path.join(DATA_DIR, "generated_docs")
                    os.makedirs(docs_dir, exist_ok=True)
                    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
                    client_slug = re.sub(
                        r"[^a-zA-Z0-9_\-]", "_", data.get("client_name") or "Client"
                    )
                    doc_filename = f"{timestamp}_{client_slug}.zip"
                    doc_path = os.path.join(docs_dir, doc_filename)
                    with open(doc_path, "wb") as df:
                        df.write(zip_bytes)
                    logger.info(
                        "Document saved: %s (%d bytes)", doc_filename, len(zip_bytes)
                    )
                    _cleanup_old_docs(docs_dir)
                except Exception as doc_err:
                    logger.warning("Could not save document copy: %s", doc_err)
                    doc_filename = None

                self.send_response(200)
                self.send_header("Content-Type", "application/zip")
                self.send_header(
                    "Content-Disposition",
                    f'attachment; filename="{client_name}_Media_Plan_Bundle.zip"',
                )
                self.send_header("Content-Length", str(len(zip_bytes)))
                self.send_header("X-Plan-Id", _plan_id)
                self.send_header("Access-Control-Expose-Headers", "X-Plan-Id")
                self.end_headers()
                self.wfile.write(zip_bytes)
                generation_time = time.time() - _gen_start
                log_request(
                    data,
                    "success",
                    file_size=len(zip_bytes),
                    generation_time=generation_time,
                    doc_filename=doc_filename,
                )
                if _metrics:
                    _metrics.record_generation(generation_time)
                # PostHog: track successful generation + funnel completion
                try:
                    from posthog_tracker import (
                        track_plan_generated,
                        track_file_upload,
                        track_plan_completed,
                    )

                    track_plan_generated(
                        email=data.get("requester_email", "anonymous"),
                        client=data.get("client_name") or "",
                        industry=data.get("industry") or "",
                        budget=str(
                            data.get("budget") or "" or data.get("budget_range") or ""
                        ),
                        roles=data.get("target_roles") or data.get("roles") or [],
                        gen_time=generation_time,
                        file_size=len(zip_bytes),
                    )
                    track_plan_completed(
                        email=data.get("requester_email", "anonymous"),
                        client=data.get("client_name") or "",
                        industry=data.get("industry") or "",
                        budget=str(
                            data.get("budget") or "" or data.get("budget_range") or ""
                        ),
                        roles_count=len(
                            data.get("target_roles") or data.get("roles") or []
                        ),
                        gen_time=generation_time,
                        file_size=len(zip_bytes),
                    )
                    # Track file uploads if any (use _upload_counts since fields were popped earlier)
                    _email = data.get("requester_email", "anonymous")
                    if _upload_counts.get("briefs") or 0:
                        track_file_upload(_email, "brief", _upload_counts["briefs"])
                    if _upload_counts.get("transcripts") or 0:
                        track_file_upload(
                            _email, "transcript", _upload_counts["transcripts"]
                        )
                    if _upload_counts.get("historical") or 0:
                        track_file_upload(
                            _email, "historical", _upload_counts["historical"]
                        )
                except Exception:
                    pass

                # ── Slack notification: sync plan generated (zip bundle) ──
                try:
                    from slack_plan_notifier import notify_plan_generated

                    _budget_alloc = data.get("_budget_allocation") or {}
                    _ch_allocs = _budget_alloc.get("channel_allocations") or {}
                    _channels_list = list(_ch_allocs.keys()) if _ch_allocs else []
                    _roles = data.get("target_roles") or data.get("roles") or []
                    _locs = data.get("locations") or []

                    notify_plan_generated(
                        {
                            "client_name": data.get("client_name") or "Unknown",
                            "industry": (
                                data.get("industry_label")
                                or data.get("industry")
                                or "N/A"
                            ),
                            "budget": data.get("budget")
                            or data.get("budget_range")
                            or "N/A",
                            "user_name": data.get("requester_name") or "Unknown",
                            "user_email": data.get("requester_email") or "Unknown",
                            "role_count": (
                                len(_roles) if isinstance(_roles, list) else 0
                            ),
                            "location_count": (
                                len(_locs) if isinstance(_locs, list) else 0
                            ),
                            "channel_count": len(_channels_list),
                            "channels_list": _channels_list,
                            "generation_time_seconds": generation_time,
                            "job_id": "",
                            "filename": f"{client_name}_Media_Plan_Bundle.zip",
                            "timestamp": time.strftime(
                                "%Y-%m-%d %H:%M:%S UTC", time.gmtime()
                            ),
                        }
                    )
                except Exception:
                    pass  # Slack notification is best-effort

            else:
                # Fallback to Excel only if PPT fails — save Excel as doc copy
                doc_filename = None
                try:
                    docs_dir = os.path.join(DATA_DIR, "generated_docs")
                    os.makedirs(docs_dir, exist_ok=True)
                    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
                    client_slug = re.sub(
                        r"[^a-zA-Z0-9_\-]", "_", data.get("client_name") or "Client"
                    )
                    # Wrap the Excel in a ZIP for consistent storage
                    doc_zip_buffer = io.BytesIO()
                    with zipfile.ZipFile(
                        doc_zip_buffer, "w", zipfile.ZIP_DEFLATED
                    ) as zf:
                        zf.writestr(f"{_descriptive_filename}.xlsx", excel_bytes)
                    doc_filename = f"{timestamp}_{client_slug}.zip"
                    doc_path = os.path.join(docs_dir, doc_filename)
                    with open(doc_path, "wb") as df:
                        df.write(doc_zip_buffer.getvalue())
                    logger.info("Document saved (Excel only): %s", doc_filename)
                except Exception as doc_err:
                    logger.warning("Could not save document copy: %s", doc_err)
                    doc_filename = None

                self.send_response(200)
                self.send_header(
                    "Content-Type",
                    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                )
                self.send_header(
                    "Content-Disposition",
                    f'attachment; filename="{_descriptive_filename}.xlsx"',
                )
                self.send_header("Content-Length", str(len(excel_bytes)))
                self.send_header("X-Plan-Id", _plan_id)
                self.send_header(
                    "Access-Control-Expose-Headers", "X-Plan-Id, X-PPT-Warning"
                )
                if pptx_warning:
                    self.send_header("X-PPT-Warning", pptx_warning)
                self.end_headers()
                self.wfile.write(excel_bytes)
                generation_time = time.time() - _gen_start
                log_request(
                    data,
                    "success",
                    file_size=len(excel_bytes),
                    generation_time=generation_time,
                    doc_filename=doc_filename,
                )
                if _metrics:
                    _metrics.record_generation(generation_time)

                # ── Slack notification: sync plan generated (Excel only) ──
                try:
                    from slack_plan_notifier import notify_plan_generated

                    _budget_alloc = data.get("_budget_allocation") or {}
                    _ch_allocs = _budget_alloc.get("channel_allocations") or {}
                    _channels_list = list(_ch_allocs.keys()) if _ch_allocs else []
                    _roles = data.get("target_roles") or data.get("roles") or []
                    _locs = data.get("locations") or []

                    notify_plan_generated(
                        {
                            "client_name": data.get("client_name") or "Unknown",
                            "industry": (
                                data.get("industry_label")
                                or data.get("industry")
                                or "N/A"
                            ),
                            "budget": data.get("budget")
                            or data.get("budget_range")
                            or "N/A",
                            "user_name": data.get("requester_name") or "Unknown",
                            "user_email": data.get("requester_email") or "Unknown",
                            "role_count": (
                                len(_roles) if isinstance(_roles, list) else 0
                            ),
                            "location_count": (
                                len(_locs) if isinstance(_locs, list) else 0
                            ),
                            "channel_count": len(_channels_list),
                            "channels_list": _channels_list,
                            "generation_time_seconds": generation_time,
                            "job_id": "",
                            "filename": f"{client_name}_Media_Plan.xlsx",
                            "timestamp": time.strftime(
                                "%Y-%m-%d %H:%M:%S UTC", time.gmtime()
                            ),
                        }
                    )
                except Exception:
                    pass  # Slack notification is best-effort

        elif path == "/api/chat":
            # ── Nova Chat Endpoint ──
            # S48: Chat auth is more permissive than /api/generate.
            # Allow: @joveo.com users, admin key, OR same-origin requests
            # (the embedded widget on media-plan-generator.onrender.com itself)
            _chat_auth_ok = (
                self._check_joveo_auth()
                or self._check_admin_auth()
                or "media-plan-generator.onrender.com"
                in (self.headers.get("Origin") or self.headers.get("Referer") or "")
                or "nova.joveo.com"
                in (self.headers.get("Origin") or self.headers.get("Referer") or "")
                or "localhost"
                in (self.headers.get("Origin") or self.headers.get("Referer") or "")
            )
            if not _chat_auth_ok:
                self._send_error(
                    "Authentication required. Please sign in with your @joveo.com account.",
                    "AUTH_REQUIRED",
                    401,
                )
                return
            if not self._check_rate_limit() or not self._check_global_chat_rate_limit():
                self._send_error(
                    "Rate limit exceeded. Please wait a moment.", "RATE_LIMITED", 429
                )
                return
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
            except (ValueError, TypeError):
                content_len = 0
            if content_len <= 0:
                self._send_error("Empty request body", "VALIDATION_ERROR", 400)
                return
            if content_len > 100 * 1024:  # 100KB limit for chat
                self._send_error("Request too large", "VALIDATION_ERROR", 413)
                return
            body = self.rfile.read(content_len)
            try:
                data = json.loads(body) if body else {}
            except (json.JSONDecodeError, RecursionError, ValueError):
                self._send_error("Invalid JSON", "VALIDATION_ERROR", 400)
                return
            # ── Type validation: reject non-string message, non-list history ──
            _raw_msg = data.get("message")
            if _raw_msg is not None and not isinstance(_raw_msg, str):
                self._send_error(
                    "Field 'message' must be a string", "VALIDATION_ERROR", 400
                )
                return
            _raw_hist = data.get("history")
            if _raw_hist is not None and not isinstance(_raw_hist, list):
                self._send_error(
                    "Field 'history' must be a list", "VALIDATION_ERROR", 400
                )
                return
            # ── Sanitize chat message to prevent XSS ──
            if "message" in data:
                data["message"] = _sanitize_chat_input(data.get("message"))
                data["message"] = _dedup_repetitive_input(data["message"])
            # ── API HARDENING: Input length limits for chat (Issue 2) ──
            _chat_msg = data.get("message") or ""
            if len(_chat_msg) > 10_000:
                self._send_error(
                    "Message exceeds 10,000 character limit", "VALIDATION_ERROR", 400
                )
                return
            _chat_history = data.get("history") or []
            # Coerce non-list history to empty list (don't crash on bad input)
            if not isinstance(_chat_history, list):
                _chat_history = []
                data["history"] = _chat_history
            # Filter out messages with null/empty content
            _chat_history = [
                m
                for m in _chat_history
                if isinstance(m, dict) and (m.get("content") or "").strip()
            ]
            data["history"] = _chat_history
            if len(_chat_history) > 50:
                self._send_error(
                    "Chat history exceeds 50 messages limit", "VALIDATION_ERROR", 400
                )
                return
            _CHAT_REQUEST_TIMEOUT: float = (
                110.0  # seconds -- S40: 90s still too tight for 6+ tool queries
                # 6 tools x 8s = 48s tools + 25s synthesis + enrichment = needs ~80-90s
                # 110s gives comfortable margin. gunicorn kill is 120s, 10s safety buffer.
                # Partial accumulator returns collected data if synthesis times out.
            )
            try:
                from nova import handle_chat_request

                _chat_span_fn = getattr(self, "_sentry_span", lambda o, n: _nullctx())

                # ── Timeout wrapper: enrichment + LLM call share budget ──
                # Enrichment (~8s) runs inside the executor so it doesn't
                # steal time from the LLM + tool loop.
                def _enriched_chat(d: dict, partial_accumulator: dict = None) -> dict:
                    import time as _t

                    _t0 = _t.time()
                    # S25: Thread outer deadline so tool loop can compute
                    # a dynamic budget that accounts for enrichment time.
                    _outer_deadline = (
                        _t0 + _CHAT_REQUEST_TIMEOUT - 5
                    )  # 5s safety margin

                    # S28: Early greeting exit -- skip enrichment for trivial queries
                    _msg_lower = (d.get("message") or "").strip().lower()
                    _msg_words = _msg_lower.split()
                    _is_trivial = len(_msg_words) <= 3 and re.match(
                        r"^(hi|hello|hey|hola|howdy|sup|yo|wassup|wazzup|whaddup|"
                        r"good morning|good afternoon|good evening|"
                        r"thanks|thank you|thx|ty|ok|okay|got it|sure|yes|no|yep|nope|"
                        r"bye|goodbye|see you|later|nice|great|cool|awesome|perfect|"
                        r"sounds good|alright|how are you|what'?s up|whats up|"
                        r"how'?s it going|how do you do|greetings|namaste)\b",
                        _msg_lower,
                    )

                    # S32: Gibberish/nonsense detection -- skip enrichment + tools
                    _is_gibberish_input = False
                    if not _is_trivial and len(_msg_words) <= 5:

                        def _is_gibberish(text: str) -> bool:
                            """Detect gibberish/keyboard mash/nonsense input."""
                            words = text.split()
                            if not words:
                                return True
                            vowels = set("aeiou")
                            real_word_count = 0
                            for w in words:
                                w_clean = re.sub(r"[^a-z]", "", w.lower())
                                if len(w_clean) < 2:
                                    continue
                                if vowels & set(w_clean):
                                    real_word_count += 1
                            # No real words (words with vowels) => gibberish
                            if real_word_count == 0 and len(text.strip()) > 0:
                                return True
                            # High consonant ratio (keyboard mash like "asdfghjkl")
                            alpha_chars = re.sub(r"[^a-z]", "", text.lower())
                            if len(alpha_chars) >= 4:
                                vowel_count = sum(1 for c in alpha_chars if c in vowels)
                                if vowel_count / len(alpha_chars) < 0.15:
                                    return True
                            return False

                        if _is_gibberish(_msg_lower):
                            _is_trivial = True
                            _is_gibberish_input = True
                            print(
                                f"[CHAT DEBUG] gibberish detected -- skipping enrichment + tools",
                                flush=True,
                            )

                    if _is_trivial and not _is_gibberish_input:
                        print(
                            f"[CHAT DEBUG] greeting detected -- skipping enrichment",
                            flush=True,
                        )
                    elif _is_gibberish_input:
                        # Return fast clarification response without LLM/tool calls
                        return {
                            "response": (
                                "I didn't quite understand that. Could you rephrase your question? "
                                "I can help with recruitment media planning, salary benchmarks, "
                                "job market trends, and more."
                            ),
                            "sources": [],
                            "confidence": 0.0,
                            "tools_used": [],
                            "llm_provider": "none (gibberish early-exit)",
                        }
                    else:
                        print(f"[CHAT DEBUG] enrichment starting", flush=True)
                        try:
                            _enrich_chat_context(d, d.get("message") or "")
                        except Exception as _enrich_err:
                            logger.error(
                                "Chat enrichment failed (non-fatal): %s",
                                _enrich_err,
                                exc_info=True,
                            )
                    _t1 = _t.time()
                    print(
                        f"[CHAT DEBUG] enrichment done in {_t1-_t0:.1f}s, calling handle_chat_request",
                        flush=True,
                    )
                    result = handle_chat_request(
                        d,
                        outer_deadline=_outer_deadline,
                        partial_accumulator=partial_accumulator,
                    )
                    _t2 = _t.time()
                    print(
                        f"[CHAT DEBUG] handle_chat_request done in {_t2-_t1:.1f}s, provider={result.get('llm_provider','?')}, tools={result.get('tools_used',[])}",
                        flush=True,
                    )
                    return result

                # S32: Thread-safe partial result accumulator -- background thread
                # writes tool results here incrementally; on timeout the main
                # thread reads whatever was collected and synthesizes a partial
                # response instead of a generic error.
                _partial_results: dict = {
                    "tools_used": [],
                    "tool_data": [],
                    "message": data.get("message") or "",
                }
                _chat_future = _background_executor.submit(
                    _enriched_chat, data, _partial_results
                )
                try:
                    with _chat_span_fn("nova.chat", "Nova chat LLM routing + response"):
                        response = _chat_future.result(timeout=_CHAT_REQUEST_TIMEOUT)
                except (TimeoutError, Exception) as _chat_timeout_err:
                    _chat_future.cancel()
                    _is_timeout = isinstance(_chat_timeout_err, TimeoutError)
                    logger.error(
                        "Chat request %s after %.0fs: %s",
                        "timed out" if _is_timeout else "failed",
                        _CHAT_REQUEST_TIMEOUT,
                        _chat_timeout_err,
                        exc_info=True,
                    )

                    # S32: Use partial results from the accumulator if available
                    _collected_tools = _partial_results.get("tool_data") or []
                    _collected_names = _partial_results.get("tools_used") or []
                    if _is_timeout and _collected_tools:
                        _partial_text = (
                            f"I gathered data from {len(_collected_names)} tool(s) "
                            f"({', '.join(_collected_names)}) before running out of time. "
                            f"Here's what I found:\n\n"
                        )
                        for _td in _collected_tools:
                            _partial_text += (
                                f"**{_td['tool']}**: {_td['result_summary'][:400]}\n\n"
                            )
                        _partial_text += (
                            "\n---\n*This is a partial response due to timeout. "
                            "Try breaking your question into smaller parts for a complete analysis.*"
                        )
                        response = {
                            "response": _partial_text,
                            "sources": [],
                            "confidence": 0.4,
                            "tools_used": _collected_names,
                            "error": "Request timed out - partial response from collected data",
                        }
                    else:
                        response = {
                            "response": (
                                "I wasn't able to complete my full analysis within the time limit, "
                                "but here's what I can tell you based on what I gathered:\n\n"
                                "Your question required multiple data lookups that took longer than expected. "
                                "Here are some suggestions:\n\n"
                                "1. **Try breaking your question into smaller parts** -- ask about one specific aspect at a time\n"
                                "2. **Be more specific** -- instead of comparing multiple cities, ask about one city first\n"
                                "3. **Try again** -- server load varies and the same question may succeed on retry\n\n"
                                "I apologize for the inconvenience. Nova processes 56+ data tools and sometimes "
                                "complex queries need a bit more time."
                                if _is_timeout
                                else "An error occurred processing your request. Please try again."
                            ),
                            "sources": [],
                            "confidence": 0.3 if _is_timeout else 0.0,
                            "tools_used": [],
                            "error": (
                                "Request timed out - partial response provided"
                                if _is_timeout
                                else str(_chat_timeout_err)[:200]
                            ),
                        }
                except BaseException as _chat_base_err:
                    # Catch SystemExit / GreenletExit from gevent worker timeout
                    _chat_future.cancel()
                    logger.error(
                        "Chat request killed (BaseException): %s",
                        _chat_base_err,
                        exc_info=True,
                    )
                    response = {
                        "response": "The server was too busy to complete your request. Please try again.",
                        "sources": [],
                        "confidence": 0.0,
                        "tools_used": [],
                        "error": "Worker interrupted",
                    }

                # ── Generate suggested follow-up questions ──
                try:
                    _resp_text = response.get("response") or ""
                    _user_msg_text = data.get("message") or ""
                    if _resp_text and _user_msg_text:
                        _followups = _generate_followup_questions(
                            _resp_text, _user_msg_text
                        )
                        if _followups:
                            response["suggested_followups"] = _followups
                except Exception as _fq_err:
                    logger.error(
                        "Follow-up generation failed: %s", _fq_err, exc_info=True
                    )

                # NOTE: Chat path routing is recorded inside nova.py (_nova_metrics.record_chat)
                # which forwards to MetricsCollector. Do NOT double-count here.
            except Exception as chat_err:
                logger.error("Nova chat error: %s", chat_err, exc_info=True)
                response = {
                    "response": "An error occurred processing your request. Please try again.",
                    "sources": [],
                    "confidence": 0.0,
                    "tools_used": [],
                    "error": "Internal error processing request",
                }
            self._send_json(response)

            # ── PostHog: Track chat message ──
            self._ph_track(
                "chat_message_sent",
                {
                    "message_length": len(data.get("message") or ""),
                    "has_attachments": bool(data.get("attachments")),
                    "provider_used": response.get("llm_provider") or "",
                    "endpoint": "/api/chat",
                },
            )

            # ── Persist conversation turn to Supabase (async, non-blocking) ──
            try:
                _conv_id = data.get("conversation_id") or ""
                _user_msg = data.get("message") or ""
                _assistant_resp = response.get("response") or ""
                _model = response.get("llm_provider") or ""
                _sources = response.get("sources") or []
                _conf = response.get("confidence") or 0.0
                if _conv_id and _user_msg:
                    submit_background_task(
                        _save_conversation_turn,
                        _conv_id,
                        _user_msg,
                        _assistant_resp,
                        _model,
                        _sources,
                        _conf,
                        task_name="supabase-persist",
                    )
            except Exception as _persist_err:
                logger.error(
                    f"Conversation persist error: {_persist_err}", exc_info=True
                )

        elif path == "/api/chat/stream":
            # ── Nova Chat SSE Streaming Endpoint ──
            # S48: Same permissive auth as /api/chat
            _stream_auth_ok = (
                self._check_joveo_auth()
                or self._check_admin_auth()
                or "media-plan-generator.onrender.com"
                in (self.headers.get("Origin") or self.headers.get("Referer") or "")
                or "nova.joveo.com"
                in (self.headers.get("Origin") or self.headers.get("Referer") or "")
                or "localhost"
                in (self.headers.get("Origin") or self.headers.get("Referer") or "")
            )
            if not _stream_auth_ok:
                self._send_error(
                    "Authentication required. Please sign in with your @joveo.com account.",
                    "AUTH_REQUIRED",
                    401,
                )
                return
            if not self._check_rate_limit() or not self._check_global_chat_rate_limit():
                self._send_error(
                    "Rate limit exceeded. Please wait a moment.", "RATE_LIMITED", 429
                )
                return
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
            except (ValueError, TypeError):
                content_len = 0
            if content_len <= 0 or content_len > 100 * 1024:
                if content_len <= 0:
                    self._send_error("Empty request body", "VALIDATION_ERROR", 400)
                else:
                    self._send_error("Request too large", "VALIDATION_ERROR", 413)
                return
            body = self.rfile.read(content_len)
            try:
                data = json.loads(body) if body else {}
            except (json.JSONDecodeError, RecursionError, ValueError):
                self._send_error("Invalid JSON", "VALIDATION_ERROR", 400)
                return
            # ── Type validation: reject non-string message, non-list history ──
            _raw_msg = data.get("message")
            if _raw_msg is not None and not isinstance(_raw_msg, str):
                self._send_error(
                    "Field 'message' must be a string", "VALIDATION_ERROR", 400
                )
                return
            _raw_hist = data.get("history")
            if _raw_hist is not None and not isinstance(_raw_hist, list):
                self._send_error(
                    "Field 'history' must be a list", "VALIDATION_ERROR", 400
                )
                return
            # ── Sanitize chat message to prevent XSS ──
            if "message" in data:
                data["message"] = _sanitize_chat_input(data.get("message"))
                data["message"] = _dedup_repetitive_input(data["message"])
            # ── API HARDENING: Input length limits for chat/stream (Issue 2) ──
            _stream_msg = data.get("message") or ""
            if len(_stream_msg) > 10_000:
                self._send_error(
                    "Message exceeds 10,000 character limit", "VALIDATION_ERROR", 400
                )
                return
            _stream_history = data.get("history") or []
            # Coerce non-list history to empty list (don't crash on bad input)
            if not isinstance(_stream_history, list):
                _stream_history = []
                data["history"] = _stream_history
            # Filter out messages with null/empty content
            _stream_history = [
                m
                for m in _stream_history
                if isinstance(m, dict) and (m.get("content") or "").strip()
            ]
            data["history"] = _stream_history
            if len(_stream_history) > 50:
                self._send_error(
                    "Chat history exceeds 50 messages limit", "VALIDATION_ERROR", 400
                )
                return

            # ── PostHog: Track chat stream message ──
            self._ph_track(
                "chat_message_sent",
                {
                    "message_length": len(data.get("message") or ""),
                    "has_attachments": bool(data.get("attachments")),
                    "provider_used": "",  # Not known yet at stream start
                    "endpoint": "/api/chat/stream",
                },
            )

            conversation_id = data.get("conversation_id") or str(uuid.uuid4())
            client_session_token = data.get("session_token") or ""

            # Issue 4: Verify session token for conversation ownership
            _conv_session_token = ""
            try:
                from nova_persistence import (
                    get_or_create_conversation,
                    verify_conversation_token,
                )

                conv_row = get_or_create_conversation(conversation_id)
                if conv_row:
                    _conv_session_token = conv_row.get("session_token") or ""
                    # If conversation is new (just created), pass back the token
                    # If existing, verify the client's token matches
                    if (
                        _conv_session_token
                        and client_session_token
                        and not verify_conversation_token(
                            conversation_id, client_session_token
                        )
                    ):
                        logger.warning(
                            "Session token mismatch for conversation %s",
                            conversation_id,
                        )
                        self._send_error(
                            "Unauthorized: invalid session token", "AUTH_ERROR", 403
                        )
                        return
            except ImportError:
                pass
            except Exception as _auth_err:
                logger.error(
                    "Conversation auth check failed: %s", _auth_err, exc_info=True
                )

            cancel_event = _register_stream(conversation_id)

            # Send SSE headers
            self.send_response(200)
            self.send_header("Content-Type", "text/event-stream")
            self.send_header("Cache-Control", "no-cache")
            self.send_header("Connection", "keep-alive")
            self.send_header("X-Accel-Buffering", "no")
            cors_origin = self._get_cors_origin()
            if cors_origin:
                self.send_header("Access-Control-Allow-Origin", cors_origin)
                self.send_header("Access-Control-Allow-Credentials", "true")
            self.end_headers()

            try:
                # ── Shared enrichment: file parsing + parallel API calls ──
                try:
                    _enrich_chat_context(data, data.get("message") or "")
                except Exception as _stream_enrich_err:
                    logger.error(
                        "Stream enrichment failed (non-fatal): %s",
                        _stream_enrich_err,
                        exc_info=True,
                    )

                # Check for cancellation before LLM call
                if cancel_event.is_set():
                    _cancel_evt = json.dumps(
                        {
                            "token": "",
                            "done": True,
                            "cancelled": True,
                            "full_response": "Your request was cancelled.",
                        }
                    )
                    self.wfile.write(f"data: {_cancel_evt}\n\n".encode())
                    self.wfile.flush()
                    return

                # ── True token-level streaming via Nova ──
                from nova import handle_chat_request_stream

                _chat_span_fn = getattr(self, "_sentry_span", lambda o, n: _nullctx())

                full_response = ""
                sources: list = []
                confidence = None  # S47: None means "don't show badge"
                tools_used: list = []
                model_used = ""
                message_id = str(uuid.uuid4())
                _token_usage: dict = {}
                # S62: Transparency panel fields forwarded from done chunk.
                _sse_fast_path = ""
                _sse_kb_files_queried: list = []
                _sse_timing_ms = None
                _last_event_time = time.time()
                _KEEPALIVE_INTERVAL = 15  # seconds between keepalive heartbeats
                _event_time_lock = (
                    threading.Lock()
                )  # protects _last_event_time across threads
                # S51 FIX: serialize concurrent writes to self.wfile between the
                # keepalive thread and the main generator thread. Without this
                # lock, a keepalive event can be interleaved mid-write with a
                # token/done event, corrupting the SSE stream so the client
                # never parses the final full_response and falls back to the
                # canned "I'm having trouble" client-side message.
                _wfile_lock = threading.Lock()

                def _sse_write(payload: str) -> bool:
                    """Write one SSE event atomically. Returns False on disconnect."""
                    try:
                        with _wfile_lock:
                            self.wfile.write(f"data: {payload}\n\n".encode())
                            self.wfile.flush()
                        return True
                    except (BrokenPipeError, ConnectionResetError, OSError):
                        return False

                # Background keepalive sender: emits {"keepalive": true} every 15s
                # while no tokens are flowing (prevents client timeout during tool-use)
                _keepalive_stop = threading.Event()

                def _send_keepalives() -> None:
                    """Send SSE keepalive heartbeats until stopped."""
                    nonlocal _last_event_time
                    while not _keepalive_stop.is_set():
                        _keepalive_stop.wait(_KEEPALIVE_INTERVAL)
                        if _keepalive_stop.is_set():
                            break
                        with _event_time_lock:
                            _elapsed = time.time() - _last_event_time
                        if _elapsed >= _KEEPALIVE_INTERVAL:
                            if not _sse_write(json.dumps({"keepalive": True})):
                                logger.debug("SSE keepalive: client disconnected")
                                break
                            with _event_time_lock:
                                _last_event_time = time.time()

                _ka_thread = threading.Thread(
                    target=_send_keepalives, daemon=True, name="sse-keepalive"
                )
                _ka_thread.start()

                # ── Immediate status event so frontend shows progress ──
                if not _sse_write(
                    json.dumps({"type": "status", "status": "Thinking..."})
                ):
                    _keepalive_stop.set()
                    return
                with _event_time_lock:
                    _last_event_time = time.time()

                try:
                    with _chat_span_fn("nova.chat.stream", "Nova SSE streaming chat"):
                        for chunk in handle_chat_request_stream(
                            data, cancel_event=cancel_event
                        ):
                            # Check cancellation on each chunk
                            if cancel_event.is_set():
                                _sse_write(
                                    json.dumps(
                                        {
                                            "token": "",
                                            "done": True,
                                            "cancelled": True,
                                            "full_response": full_response,
                                            "sources": sources,
                                        }
                                    )
                                )
                                return

                            if chunk.get("done"):
                                # Final metadata chunk -- extract info
                                full_response = (
                                    chunk.get("full_response") or full_response
                                )
                                sources = chunk.get("sources") or sources
                                # S47: Preserve None confidence (means "don't show badge")
                                _chunk_conf = chunk.get("confidence")
                                if _chunk_conf is not None:
                                    confidence = _chunk_conf
                                tools_used = chunk.get("tools_used") or tools_used
                                model_used = chunk.get("llm_provider") or model_used
                                _token_usage = chunk.get("token_usage") or {}
                                # S62: Pull transparency panel fields from the
                                # Nova done chunk so the final SSE event can
                                # forward them to the client. Without this the
                                # panel saw fast_path=None and kb_files_queried=[]
                                # even when the fast-path had actually fired.
                                _sse_fast_path = chunk.get("fast_path") or ""
                                _sse_kb_files_queried = (
                                    chunk.get("kb_files_queried") or []
                                )
                                _sse_timing_ms = chunk.get("timing_ms")
                            elif chunk.get("type") in ("tool_start", "tool_complete"):
                                # S18: Real-time tool status -- forward to client
                                if not _sse_write(
                                    json.dumps(
                                        {
                                            "type": chunk["type"],
                                            "tool": chunk.get("tool", ""),
                                            "label": chunk.get("label", ""),
                                            "done": False,
                                        }
                                    )
                                ):
                                    break
                                with _event_time_lock:
                                    _last_event_time = time.time()
                            elif chunk.get("status"):
                                # Status/progress event -- forward to client
                                if not _sse_write(
                                    json.dumps(
                                        {"status": chunk["status"], "done": False}
                                    )
                                ):
                                    break
                                with _event_time_lock:
                                    _last_event_time = time.time()
                            else:
                                # Token chunk -- stream to client immediately
                                token = chunk.get("token") or ""
                                if token:
                                    full_response += token
                                    if not _sse_write(
                                        json.dumps({"token": token, "done": False})
                                    ):
                                        break
                                    with _event_time_lock:
                                        _last_event_time = time.time()
                finally:
                    _keepalive_stop.set()

                # ── Generate follow-up suggestions ──
                _stream_followups: list[str] = []
                try:
                    _stream_user_msg = data.get("message") or ""
                    if full_response and _stream_user_msg:
                        _stream_followups = _generate_followup_questions(
                            full_response, _stream_user_msg
                        )
                except Exception as _sfq:
                    logger.error("Stream follow-up gen failed: %s", _sfq, exc_info=True)

                # ── Final event with metadata ──
                # S51 FIX: If the pipeline completed but produced no text, emit
                # a concrete server-side fallback instead of empty, so the client
                # widget does NOT fall back to the generic "I'm having trouble"
                # canned message. The widget only shows that message when both
                # full_response and accumulated tokens are empty.
                if not (full_response or "").strip():
                    full_response = (
                        "I couldn't complete this query in time. "
                        "This can happen on complex multi-constraint questions. "
                        "Try breaking it into steps -- for example: "
                        '"List healthcare job boards in the US" then ask a '
                        "follow-up for supply partners."
                    )
                final = json.dumps(
                    {
                        "token": "",
                        "done": True,
                        "full_response": full_response,
                        "sources": sources,
                        "confidence": confidence,
                        "tools_used": tools_used,
                        "message_id": message_id,
                        "conversation_id": conversation_id,
                        "suggested_followups": _stream_followups,
                        "session_token": _conv_session_token,
                        "token_usage": _token_usage,
                        # S62: Transparency panel payload -- forwarded from
                        # Nova stream done chunk so the "Why this answer?"
                        # footer shows the real fast-path marker and KB files
                        # that were consulted (empty list is the honest answer
                        # when the fast-path used hardcoded dicts).
                        "fast_path": _sse_fast_path,
                        "kb_files_queried": _sse_kb_files_queried,
                        "timing_ms": _sse_timing_ms,
                    }
                )
                _sse_write(final)

                # Store in idempotency cache so retries get a clean JSON response
                try:
                    _idem_body = json.dumps(
                        {
                            "success": True,
                            "error": None,
                            "code": None,
                            "data": {
                                "full_response": full_response,
                                "sources": sources,
                                "confidence": confidence,
                                "message_id": message_id,
                                "conversation_id": conversation_id,
                                "replayed": True,
                            },
                        }
                    ).encode("utf-8")
                    self._store_idempotency("/api/chat/stream", _idem_body, 200)
                except Exception:
                    pass  # Non-critical

            except (BrokenPipeError, ConnectionResetError):
                logger.info(
                    "SSE client disconnected for conversation %s", conversation_id
                )
            except Exception as stream_err:
                logger.error("SSE streaming error: %s", stream_err, exc_info=True)
                _sse_write(
                    json.dumps(
                        {
                            "token": "",
                            "done": True,
                            "full_response": (
                                "I hit an internal error processing your request. "
                                "Please try again in a moment."
                            ),
                            "error": "An error occurred processing your request.",
                        }
                    )
                )
            finally:
                # Persist to Supabase asynchronously -- in finally so
                # BrokenPipeError / ConnectionResetError doesn't skip it
                _user_msg = data.get("message") or ""
                if conversation_id and _user_msg and full_response:
                    threading.Thread(
                        target=_save_conversation_turn,
                        args=(
                            conversation_id,
                            _user_msg,
                            full_response,
                            model_used,
                            sources,
                            confidence,
                        ),
                        daemon=True,
                        name="supabase-stream-persist",
                    ).start()
                _unregister_stream(conversation_id)

        elif path == "/api/chat/stop":
            # ── Stop/Cancel an active streaming request ──
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
            except (ValueError, TypeError):
                content_len = 0
            body = self.rfile.read(content_len) if content_len > 0 else b""
            try:
                data = json.loads(body) if body else {}
            except json.JSONDecodeError:
                data = {}
            conv_id = data.get("conversation_id") or ""
            cancelled = _cancel_stream(conv_id) if conv_id else False
            resp_body = json.dumps(
                {
                    "cancelled": cancelled,
                    "conversation_id": conv_id,
                }
            ).encode("utf-8")
            self.send_response(200)
            self.send_header("Content-Type", "application/json")
            cors_origin = self._get_cors_origin()
            if cors_origin:
                self.send_header("Access-Control-Allow-Origin", cors_origin)
            self.send_header("Content-Length", str(len(resp_body)))
            self.end_headers()
            self.wfile.write(resp_body)

        elif path == "/api/nova/summarize":
            # ── Conversation Summarization Endpoint ──
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
            except (ValueError, TypeError):
                content_len = 0
            if content_len <= 0 or content_len > 100 * 1024:
                self._send_error("Invalid request body", "VALIDATION_ERROR", 400)
                return
            body = self.rfile.read(content_len)
            try:
                data = json.loads(body) if body else {}
            except (json.JSONDecodeError, ValueError):
                self._send_error("Invalid JSON", "VALIDATION_ERROR", 400)
                return

            _sum_conv_id = data.get("conversation_id") or ""
            if not _sum_conv_id or not isinstance(_sum_conv_id, str):
                self._send_error("conversation_id is required", "VALIDATION_ERROR", 400)
                return

            # Optional: pass messages directly (for client-side conversations)
            _sum_messages = data.get("messages")
            if _sum_messages and not isinstance(_sum_messages, list):
                self._send_error("messages must be an array", "VALIDATION_ERROR", 400)
                return

            try:
                from nova import summarize_conversation

                _sum_result = summarize_conversation(
                    conversation_id=_sum_conv_id,
                    messages=_sum_messages,
                )
                self._send_json(_sum_result)
            except Exception as _sum_err:
                logger.error("Summarize endpoint error: %s", _sum_err, exc_info=True)
                self._send_json(
                    {
                        "error": "Failed to generate summary",
                        "conversation_id": _sum_conv_id,
                    },
                    status_code=500,
                )

        # NOTE: /api/campaign/save now handled by routes/campaign.py
        elif path == "/api/chat/feedback":
            # ── Chat message feedback (thumbs up/down) ──
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
            except (ValueError, TypeError):
                content_len = 0
            body = self.rfile.read(content_len) if content_len > 0 else b""
            try:
                data = json.loads(body) if body else {}
            except json.JSONDecodeError:
                data = {}
            msg_id = data.get("message_id") or ""
            rating = data.get("rating") or ""
            conv_id = data.get("conversation_id") or ""
            feedback_text = data.get("feedback_text") or ""
            # ── API HARDENING: feedback_text length limit ──
            if len(feedback_text) > 5_000:
                self._send_error(
                    "feedback_text exceeds 5,000 character limit",
                    "VALIDATION_ERROR",
                    400,
                )
                return
            if not msg_id or rating not in ("up", "down"):
                resp_body = json.dumps(
                    {
                        "error": "Required: message_id (string), rating ('up' or 'down')",
                    }
                ).encode("utf-8")
                self.send_response(400)
                self.send_header("Content-Type", "application/json")
                cors_origin = self._get_cors_origin()
                if cors_origin:
                    self.send_header("Access-Control-Allow-Origin", cors_origin)
                self.send_header("Content-Length", str(len(resp_body)))
                self.end_headers()
                self.wfile.write(resp_body)
                return
            _store_feedback(msg_id, rating, conv_id, feedback_text)
            resp_body = json.dumps(
                {
                    "status": "ok",
                    "message_id": msg_id,
                    "rating": rating,
                }
            ).encode("utf-8")
            self.send_response(200)
            self.send_header("Content-Type", "application/json")
            cors_origin = self._get_cors_origin()
            if cors_origin:
                self.send_header("Access-Control-Allow-Origin", cors_origin)
            self.send_header("Content-Length", str(len(resp_body)))
            self.end_headers()
            self.wfile.write(resp_body)

        elif path == "/api/chat/title":
            # ── S48: Update conversation title (client-side smart title sync) ──
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
            except (ValueError, TypeError):
                content_len = 0
            body = self.rfile.read(content_len) if content_len > 0 else b""
            try:
                data = json.loads(body) if body else {}
            except json.JSONDecodeError:
                data = {}
            conv_id = (data.get("conversation_id") or "").strip()
            new_title = (data.get("title") or "").strip()
            if not conv_id or not new_title:
                self._send_error(
                    "Required: conversation_id and title",
                    "VALIDATION_ERROR",
                    400,
                )
                return
            # Sanitize title (max 100 chars, strip HTML)
            new_title = re.sub(r"<[^>]+>", "", new_title)[:100]
            try:
                from nova_persistence import update_conversation

                update_conversation(conv_id, title=new_title)
            except ImportError:
                logger.debug("nova_persistence not available for title update")
            except Exception as exc:
                logger.error(
                    "Failed to update title for %s: %s",
                    conv_id,
                    exc,
                    exc_info=True,
                )
            self._send_json(
                {"status": "ok", "conversation_id": conv_id, "title": new_title}
            )

        elif path == "/api/chat/share":
            # ── Share conversation ──
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
            except (ValueError, TypeError):
                content_len = 0
            if content_len <= 0:
                self._send_error("Empty request body", "VALIDATION_ERROR", 400)
                return
            if content_len > 200 * 1024:  # 200KB limit for share payloads
                self._send_error("Request too large", "VALIDATION_ERROR", 413)
                return
            body = self.rfile.read(content_len)
            try:
                data = json.loads(body) if body else {}
            except (json.JSONDecodeError, RecursionError, ValueError):
                self._send_error("Invalid JSON", "VALIDATION_ERROR", 400)
                return
            # ── Validate messages array ──
            _share_messages = data.get("messages")
            if not isinstance(_share_messages, list) or len(_share_messages) < 1:
                self._send_error(
                    "Payload must include 'messages' array with at least 1 message",
                    "VALIDATION_ERROR",
                    400,
                )
                return
            if len(_share_messages) > 100:
                self._send_error(
                    "Share payload exceeds 100 messages limit",
                    "VALIDATION_ERROR",
                    400,
                )
                return
            # ── Sanitize message content for XSS ──
            _has_content = False
            for _smsg in _share_messages:
                if not isinstance(_smsg, dict):
                    continue
                _smsg_content = _smsg.get("content")
                if isinstance(_smsg_content, str) and _smsg_content.strip():
                    _has_content = True
                    _smsg["content"] = _sanitize_chat_input(_smsg_content)
                # Sanitize role field too (should only be user/assistant/system)
                _smsg_role = _smsg.get("role") or ""
                if _smsg_role not in ("user", "assistant", "system"):
                    _smsg["role"] = "user"
            if not _has_content:
                self._send_error(
                    "Messages must contain at least one message with non-empty content",
                    "VALIDATION_ERROR",
                    400,
                )
                return
            # Sanitize title
            _share_title: str = _sanitize_chat_input(
                data.get("title") or "Shared Conversation"
            )
            conv_id = data.get("conversation_id") or str(uuid.uuid4())
            share_id = str(uuid.uuid4())[:8]
            # Store to Supabase (table may not exist -- graceful degrade)
            _share_saved = False
            try:
                _supabase_rest(
                    "nova_shared_conversations",
                    method="POST",
                    payload={
                        "share_id": share_id,
                        "conversation_id": conv_id,
                        "title": _share_title,
                        "messages": json.dumps(_share_messages),
                        "created_at": datetime.datetime.utcnow().isoformat() + "Z",
                    },
                )
                _share_saved = True
            except Exception as _se:
                logger.error(
                    "Share conversation save failed (table nova_shared_conversations may not exist): %s",
                    _se,
                    exc_info=True,
                )
            if not _share_saved:
                resp_body = json.dumps(
                    {
                        "status": "error",
                        "message": "Sharing is temporarily unavailable.",
                    }
                ).encode("utf-8")
                self.send_response(503)
                self.send_header("Content-Type", "application/json")
                cors_origin = self._get_cors_origin()
                if cors_origin:
                    self.send_header("Access-Control-Allow-Origin", cors_origin)
                self.send_header("Content-Length", str(len(resp_body)))
                self.end_headers()
                self.wfile.write(resp_body)
                return
            share_url = f"{self.headers.get('Origin') or ''}/nova/shared/{share_id}"
            resp_body = json.dumps(
                {"status": "ok", "share_id": share_id, "share_url": share_url}
            ).encode("utf-8")
            self.send_response(200)
            self.send_header("Content-Type", "application/json")
            cors_origin = self._get_cors_origin()
            if cors_origin:
                self.send_header("Access-Control-Allow-Origin", cors_origin)
            self.send_header("Content-Length", str(len(resp_body)))
            self.end_headers()
            self.wfile.write(resp_body)

        elif path == "/api/slack/events":
            # ── Nova Slack Event Webhook ──
            # FIX: Handle Slack retries — return 200 immediately to prevent duplicates
            retry_num = self.headers.get("X-Slack-Retry-Num")
            retry_reason = self.headers.get("X-Slack-Retry-Reason")
            if retry_num:
                logger.info(
                    "Slack retry #%s (reason: %s) — acknowledging without reprocessing",
                    retry_num,
                    retry_reason,
                )
                self.send_response(200)
                self.send_header("Content-Type", "application/json")
                self.send_header("X-Slack-No-Retry", "1")
                self.end_headers()
                self.wfile.write(json.dumps({"ok": True}).encode())
                return
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
            except (ValueError, TypeError):
                content_len = 0
            if content_len > 1 * 1024 * 1024:  # 1MB limit for Slack events
                self._send_error("Request too large", "VALIDATION_ERROR", 413)
                return
            body = self.rfile.read(content_len)
            try:
                data = json.loads(body) if body else {}
            except json.JSONDecodeError:
                self._send_error("Invalid JSON", "VALIDATION_ERROR", 400)
                return

            # Verify Slack request signature -- SECURITY: never skip on error
            # in production. Only skip if SLACK_SIGNING_SECRET is not configured
            # (i.e., the Slack bot is intentionally disabled).
            _slack_signing_secret = os.environ.get("SLACK_SIGNING_SECRET") or ""
            if _slack_signing_secret:
                try:
                    from nova_slack import get_nova_bot

                    _nova_bot = get_nova_bot()
                    _slack_ts = self.headers.get("X-Slack-Request-Timestamp") or ""
                    _slack_sig = self.headers.get("X-Slack-Signature") or ""
                    if not _nova_bot.verify_slack_signature(
                        _slack_ts,
                        body.decode("utf-8") if isinstance(body, bytes) else body,
                        _slack_sig,
                    ):
                        logger.warning("Slack signature verification failed")
                        self.send_response(403)
                        self.send_header("Content-Type", "application/json")
                        self.end_headers()
                        self.wfile.write(
                            _error_response("Invalid signature", "AUTH_REQUIRED", 401)[
                                0
                            ]
                        )
                        return
                except ImportError:
                    # nova_slack module not available -- reject the request
                    logger.error(
                        "Slack signing secret configured but nova_slack module unavailable"
                    )
                    self.send_response(500)
                    self.send_header("Content-Type", "application/json")
                    self.end_headers()
                    self.wfile.write(
                        _error_response(
                            "Slack integration unavailable", "SERVER_ERROR", 503
                        )[0]
                    )
                    return
                except Exception as sig_err:
                    # Signature verification failed due to unexpected error -- reject
                    logger.error("Slack signature verification error: %s", sig_err)
                    self.send_response(403)
                    self.send_header("Content-Type", "application/json")
                    self.end_headers()
                    self.wfile.write(
                        _error_response(
                            "Signature verification failed", "AUTH_REQUIRED", 401
                        )[0]
                    )
                    return
            else:
                logger.warning(
                    "SLACK_SIGNING_SECRET not set -- rejecting Slack event (configure to enable)"
                )
                self.send_response(403)
                self.send_header("Content-Type", "application/json")
                self.end_headers()
                self.wfile.write(
                    json.dumps(
                        {"error": "Slack signing secret not configured"}
                    ).encode()
                )
                return

            # Respond to Slack immediately (within 3 seconds), process async
            self.send_response(200)
            self.send_header("Content-Type", "application/json")
            self.send_header("X-Slack-No-Retry", "1")
            self.end_headers()

            # Handle URL verification challenge synchronously (Slack requires immediate response)
            if data.get("type") == "url_verification":
                challenge_resp = json.dumps(
                    {"challenge": data.get("challenge") or ""}
                ).encode()
                self.wfile.write(challenge_resp)
                return

            self.wfile.write(json.dumps({"ok": True}).encode())

            # Process event asynchronously via bounded thread pool (prevents thread explosion)
            def _process_slack_event_async(event_data):
                try:
                    from nova_slack import handle_slack_event

                    handle_slack_event(event_data)
                    if _metrics:
                        _metrics.record_slack_event()
                except Exception as slack_err:
                    logger.error(
                        "Nova Slack async event error: %s", slack_err, exc_info=True
                    )

            try:
                _slack_executor.submit(_process_slack_event_async, data)
            except RuntimeError:
                # Thread pool shut down or full -- process synchronously as fallback
                logger.warning(
                    "Slack thread pool unavailable, processing synchronously"
                )
                _process_slack_event_async(data)
        elif path == "/api/admin/nova":
            # ── Nova Admin API (unanswered questions management) ──
            if not self._check_admin_auth():
                self._send_error("Unauthorized", "AUTH_REQUIRED", 401)
                return
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
            except (ValueError, TypeError):
                content_len = 0
            if content_len > 1 * 1024 * 1024:  # 1MB limit for admin
                self._send_error("Request too large", "VALIDATION_ERROR", 413)
                return
            body = self.rfile.read(content_len)
            try:
                data = json.loads(body) if body else {}
            except json.JSONDecodeError:
                self._send_error("Invalid JSON", "VALIDATION_ERROR", 400)
                return
            try:
                from nova_slack import handle_admin_unanswered

                response = handle_admin_unanswered(data)
            except Exception as admin_err:
                logger.error("Nova admin error: %s", admin_err, exc_info=True)
                response = {"error": "Internal error processing request"}
            resp_body = json.dumps(response).encode("utf-8")
            self.send_response(200)
            self.send_header("Content-Type", "application/json")
            self.send_header("Content-Length", str(len(resp_body)))
            self.end_headers()
            self.wfile.write(resp_body)
        elif path == "/api/nova/chat":
            # ── Nova Standalone Chat (admin testing) ──
            if not self._check_admin_auth():
                self._send_error("Unauthorized", "AUTH_REQUIRED", 401)
                return
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
            except (ValueError, TypeError):
                content_len = 0
            if content_len > 1 * 1024 * 1024:  # 1MB limit
                self._send_error("Request too large", "VALIDATION_ERROR", 413)
                return
            body = self.rfile.read(content_len)
            try:
                data = json.loads(body) if body else {}
            except json.JSONDecodeError:
                self._send_error("Invalid JSON", "VALIDATION_ERROR", 400)
                return
            # ── Sanitize chat message to prevent XSS ──
            if "message" in data:
                data["message"] = _sanitize_chat_input(data.get("message"))
                data["message"] = _dedup_repetitive_input(data["message"])
            try:
                from nova_slack import handle_chat_standalone

                _admin_span_fn = getattr(self, "_sentry_span", lambda o, n: _nullctx())
                with _admin_span_fn("nova.admin_chat", "Nova admin chat (standalone)"):
                    response = handle_chat_standalone(data)
            except Exception as chat_err:
                logger.error("Nova chat error: %s", chat_err, exc_info=True)
                response = {"error": "Internal error processing request"}
            resp_body = json.dumps(response).encode("utf-8")
            self.send_response(200)
            self.send_header("Content-Type", "application/json")
            self.send_header("Content-Length", str(len(resp_body)))
            self.end_headers()
            self.wfile.write(resp_body)
        # ── Feature 3c: Admin API key management ──
        elif path == "/api/admin/keys":
            if not self._check_admin_auth():
                self._send_error("Unauthorized", "AUTH_REQUIRED", 401)
                return
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
            except (ValueError, TypeError):
                content_len = 0
            body = self.rfile.read(content_len) if content_len > 0 else b"{}"
            try:
                data = json.loads(body)
            except json.JSONDecodeError:
                self._send_error("Invalid JSON", "VALIDATION_ERROR", 400)
                return
            action = data.get("action", "list")
            if action == "create":
                tier = data.get("tier", "free")
                if tier not in API_KEY_TIERS:
                    self._send_json(
                        {
                            "error": f"Invalid tier. Must be one of: {list(API_KEY_TIERS.keys())}"
                        }
                    )
                    return
                label = data.get("label") or ""
                new_key = uuid.uuid4().hex
                with _api_keys_lock:
                    _api_keys_store[new_key] = {
                        "tier": tier,
                        "label": label,
                        "created": datetime.datetime.now(
                            datetime.timezone.utc
                        ).isoformat(),
                        "revoked": False,
                        "usage_minute": [],
                        "usage_day": [],
                    }
                self._send_json(
                    {
                        "action": "created",
                        "key": new_key,
                        "tier": tier,
                        "label": label,
                        "limits": API_KEY_TIERS[tier],
                    }
                )
            elif action == "revoke":
                key_to_revoke = data.get("key") or ""
                with _api_keys_lock:
                    if key_to_revoke in _api_keys_store:
                        _api_keys_store[key_to_revoke]["revoked"] = True
                        self._send_json(
                            {"action": "revoked", "key": key_to_revoke[:8] + "..."}
                        )
                    else:
                        self._send_json({"error": "Key not found"})
            elif action == "list":
                keys_list = []
                with _api_keys_lock:
                    for k, v in _api_keys_store.items():
                        keys_list.append(
                            {
                                "key_prefix": k[:8] + "...",
                                "tier": v.get("tier", "free"),
                                "label": v.get("label") or "",
                                "created": v.get("created") or "",
                                "revoked": v.get("revoked", False),
                            }
                        )
                self._send_json({"keys": keys_list, "total": len(keys_list)})
            else:
                self._send_json(
                    {"error": f"Unknown action '{action}'. Use: create, list, revoke"}
                )
        # ── Knowledge Base Backup (admin) ──
        elif path == "/api/admin/backup":
            if not self._check_admin_auth():
                self._send_error("Unauthorized", "AUTH_REQUIRED", 401)
                return
            try:
                from scripts.backup_kb import backup_knowledge_base

                backup_path = backup_knowledge_base(DATA_DIR)
                backup_file = Path(backup_path)
                self._send_json(
                    {
                        "status": "ok",
                        "filename": backup_file.name,
                        "size_bytes": backup_file.stat().st_size,
                        "path": str(backup_file),
                    }
                )
            except Exception as bk_err:
                logger.error("Backup failed: %s", bk_err, exc_info=True)
                self.send_response(500)
                self.send_header("Content-Type", "application/json")
                self.end_headers()
                self.wfile.write(
                    _error_response(f"Backup failed: {bk_err}", "SERVER_ERROR", 500)[0]
                )
        # ── Campaign Performance Tracker APIs ──
        elif path == "/api/tracker/analyze":
            if not self._check_rate_limit():
                self._send_error("Rate limit exceeded", "RATE_LIMITED", 429)
                return
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
            except (ValueError, TypeError):
                content_len = 0
            if content_len <= 0 or content_len > 50 * 1024 * 1024:
                _sz_status = 400 if content_len <= 0 else 413
                _sz_msg = (
                    "Empty request body" if content_len <= 0 else "Request too large"
                )
                self._send_error(_sz_msg, "VALIDATION_ERROR", _sz_status)
                return
            body = self.rfile.read(content_len)
            try:
                content_type = self.headers.get("Content-Type") or ""
                if "multipart/form-data" in content_type:
                    import cgi

                    # io is imported at module level (line 6)

                    environ = {
                        "REQUEST_METHOD": "POST",
                        "CONTENT_TYPE": content_type,
                        "CONTENT_LENGTH": str(content_len),
                    }
                    fs = cgi.FieldStorage(
                        fp=io.BytesIO(body), environ=environ, keep_blank_values=True
                    )
                    file_item = fs["file"] if "file" in fs else None
                    file_bytes = file_item.file.read() if file_item else b""
                    filename = file_item.filename if file_item else "upload.csv"
                    campaign_name = fs.getfirst("campaign_name", "Campaign")
                    industry = fs.getfirst("industry", "Technology")
                    roles = fs.getfirst("roles", "")
                    locations = fs.getfirst("locations", "")
                else:
                    data = json.loads(body)
                    file_bytes = None
                    filename = ""
                    campaign_name = data.get("campaign_name", "Campaign")
                    industry = data.get("industry", "Technology")
                    roles = data.get("roles") or ""
                    locations = data.get("locations") or ""
                from performance_tracker import analyze_campaign

                result = analyze_campaign(
                    file_bytes=file_bytes,
                    filename=filename,
                    campaign_name=campaign_name,
                    industry=industry,
                    roles=(
                        [r.strip() for r in roles.split(",") if r.strip()]
                        if isinstance(roles, str)
                        else roles
                    ),
                    locations=(
                        [l.strip() for l in locations.split(",") if l.strip()]
                        if isinstance(locations, str)
                        else locations
                    ),
                )
                # ── API enrichment: Adzuna salary benchmarks for CPC/CPA context ──
                if isinstance(result, dict) and _api_integrations_available:
                    tracker_benchmarks: dict = {}
                    try:
                        roles_list = (
                            [r.strip() for r in roles.split(",") if r.strip()]
                            if isinstance(roles, str)
                            else roles or []
                        )
                        primary_role = roles_list[0] if roles_list else ""
                        if primary_role and _api_adzuna:
                            salary_hist = _api_adzuna.get_salary_histogram(primary_role)
                            if salary_hist:
                                tracker_benchmarks["salary_benchmark"] = salary_hist
                            top_companies = _api_adzuna.get_top_companies(primary_role)
                            if top_companies:
                                tracker_benchmarks["top_hiring_companies"] = (
                                    top_companies
                                )
                    except Exception as api_err:
                        logger.error(
                            f"Tracker Adzuna enrichment error: {api_err}", exc_info=True
                        )
                    if tracker_benchmarks:
                        result["industry_benchmarks"] = tracker_benchmarks
                # LLM insights for Performance Tracker
                if isinstance(result, dict):
                    result["ai_insights"] = _generate_product_insights(
                        "Campaign Performance Tracker",
                        {
                            k: result.get(k)
                            for k in (
                                "summary",
                                "channel_performance",
                                "metrics",
                                "alerts",
                            )
                            if result.get(k) is not None
                        },
                        context=f"Campaign: {campaign_name}, Industry: {industry}",
                    )
                resp_body = json.dumps(result).encode("utf-8")
                self.send_response(200)
                self.send_header("Content-Type", "application/json")
                cors_origin = self._get_cors_origin()
                if cors_origin:
                    self.send_header("Access-Control-Allow-Origin", cors_origin)
                self.send_header("Content-Length", str(len(resp_body)))
                self.end_headers()
                self.wfile.write(resp_body)
            except Exception as e:
                logger.error("Tracker analyze error: %s", e, exc_info=True)
                self.send_response(500)
                self.send_header("Content-Type", "application/json")
                self.end_headers()
                self.wfile.write(
                    _error_response("Analysis failed", "SERVER_ERROR", 500)[0]
                )
        elif path == "/api/tracker/download/excel":
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
                body = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body)
                from performance_tracker import generate_performance_excel

                excel_bytes = generate_performance_excel(
                    data.get("report_data", data), data.get("client_name", "Campaign")
                )
                self.send_response(200)
                self.send_header(
                    "Content-Type",
                    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                )
                self.send_header(
                    "Content-Disposition",
                    "attachment; filename=performance_report.xlsx",
                )
                self.send_header("Content-Length", str(len(excel_bytes)))
                self.end_headers()
                self.wfile.write(excel_bytes)
            except Exception as e:
                logger.error("Tracker Excel download error: %s", e, exc_info=True)
                self.send_response(500)
                self.send_header("Content-Type", "application/json")
                self.end_headers()
                self.wfile.write(
                    _error_response("Excel generation failed", "SERVER_ERROR", 500)[0]
                )
        elif path == "/api/tracker/download/ppt":
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
                body = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body)
                from performance_tracker import generate_performance_ppt

                ppt_bytes = generate_performance_ppt(
                    data.get("report_data", data), data.get("client_name", "Campaign")
                )
                self.send_response(200)
                self.send_header(
                    "Content-Type",
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
                self.send_header(
                    "Content-Disposition",
                    "attachment; filename=performance_report.pptx",
                )
                self.send_header("Content-Length", str(len(ppt_bytes)))
                self.end_headers()
                self.wfile.write(ppt_bytes)
            except Exception as e:
                logger.error("Tracker PPT download error: %s", e, exc_info=True)
                self.send_response(500)
                self.send_header("Content-Type", "application/json")
                self.end_headers()
                self.wfile.write(
                    _error_response("PPT generation failed", "SERVER_ERROR", 500)[0]
                )
        # ── Budget Simulator APIs ──
        elif path == "/api/simulator/defaults":
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
                body = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body)
                from budget_simulator import get_default_allocation

                _sim_roles = data.get("roles", ["Software Engineer"])
                _sim_locs = data.get("locations", ["United States"])
                result = get_default_allocation(
                    total_budget=float(data.get("total_budget", 100000)),
                    industry=data.get("industry", "Technology"),
                    roles=(
                        ", ".join(_sim_roles)
                        if isinstance(_sim_roles, list)
                        else str(_sim_roles)
                    ),
                    locations=(
                        ", ".join(_sim_locs)
                        if isinstance(_sim_locs, list)
                        else str(_sim_locs)
                    ),
                )
                self._send_json(result)
            except Exception as e:
                logger.error("Simulator defaults POST error: %s", e, exc_info=True)
                self._send_json({"error": _internal_error_msg()})
        elif path == "/api/simulator/simulate":
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
                body = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body)
                from budget_simulator import simulate_scenario

                _sim_roles = data.get("roles", ["Software Engineer"])
                _sim_locs = data.get("locations", ["United States"])
                result = simulate_scenario(
                    channel_allocations=data.get(
                        "channel_allocations", data.get("allocations", {})
                    ),
                    total_budget=float(data.get("total_budget", 100000)),
                    industry=data.get("industry", "Technology"),
                    roles=(
                        ", ".join(_sim_roles)
                        if isinstance(_sim_roles, list)
                        else str(_sim_roles)
                    ),
                    locations=(
                        ", ".join(_sim_locs)
                        if isinstance(_sim_locs, list)
                        else str(_sim_locs)
                    ),
                )
                # ── API enrichment: FRED + BEA for economic context ──
                if isinstance(result, dict) and _api_integrations_available:
                    sim_api_context: dict = {}
                    try:
                        if _api_fred:
                            unemp = _api_fred.get_unemployment_rate()
                            if unemp:
                                sim_api_context["unemployment_rate"] = unemp
                            cpi = _api_fred.get_cpi_data(months=3)
                            if cpi:
                                sim_api_context["cpi_recent"] = cpi
                    except Exception as api_err:
                        logger.error(
                            f"Simulator FRED enrichment error: {api_err}", exc_info=True
                        )
                    try:
                        if _api_bea:
                            # Extract state from locations for GDP context
                            _loc_str = (
                                _sim_locs[0]
                                if isinstance(_sim_locs, list) and _sim_locs
                                else str(_sim_locs)
                            )
                            gdp = _api_bea.get_gdp_by_state()
                            if gdp:
                                sim_api_context["gdp_by_state"] = gdp
                    except Exception as api_err:
                        logger.error(
                            f"Simulator BEA enrichment error: {api_err}", exc_info=True
                        )
                    if sim_api_context:
                        result["economic_context"] = sim_api_context
                # LLM insights for Budget Simulator
                if isinstance(result, dict):
                    result["ai_insights"] = _generate_product_insights(
                        "Budget Simulator",
                        {
                            k: result.get(k)
                            for k in (
                                "summary",
                                "total_budget",
                                "channel_allocations",
                                "projected_metrics",
                            )
                            if result.get(k) is not None
                        },
                        context=f"Industry: {data.get('industry', 'Technology')}, "
                        f"Roles: {_sim_roles}, Locations: {_sim_locs}",
                    )
                self._send_json(result)
            except Exception as e:
                logger.error("Simulator simulate error: %s", e, exc_info=True)
                self._send_json({"error": _internal_error_msg()})
        elif path == "/api/simulator/optimize":
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
                body = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body)
                from budget_simulator import optimize_for_goal

                _opt_roles = data.get("roles", ["Software Engineer"])
                _opt_locs = data.get("locations", ["United States"])
                result = optimize_for_goal(
                    total_budget=float(data.get("total_budget", 100000)),
                    goal=data.get("goal", "balanced"),
                    industry=data.get("industry", "Technology"),
                    roles=(
                        ", ".join(_opt_roles)
                        if isinstance(_opt_roles, list)
                        else str(_opt_roles)
                    ),
                    locations=(
                        ", ".join(_opt_locs)
                        if isinstance(_opt_locs, list)
                        else str(_opt_locs)
                    ),
                )
                # LLM insights for Budget Optimizer
                if isinstance(result, dict):
                    result["ai_insights"] = _generate_product_insights(
                        "Budget Optimizer",
                        {
                            k: result.get(k)
                            for k in (
                                "goal",
                                "optimized_allocations",
                                "projected_metrics",
                            )
                            if result.get(k) is not None
                        },
                        context=f"Goal: {data.get('goal', 'balanced')}, "
                        f"Budget: ${data.get('total_budget', 100000):,.0f}",
                    )
                self._send_json(result)
            except Exception as e:
                logger.error("Simulator optimize error: %s", e, exc_info=True)
                self._send_json({"error": _internal_error_msg()})
        elif path == "/api/simulator/compare":
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
                body = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body)
                from budget_simulator import compare_scenarios, simulate_scenario

                raw_scenarios = data.get("scenarios") or []
                # Pre-simulate scenarios if they are raw configs (no "summary" key)
                simulated = []
                for sc in raw_scenarios:
                    if "summary" in sc:
                        simulated.append(sc)
                    else:
                        _sc_roles = sc.get("roles", ["Software Engineer"])
                        _sc_locs = sc.get("locations", ["United States"])
                        sim_result = simulate_scenario(
                            channel_allocations=sc.get(
                                "channel_allocations", sc.get("allocations", {})
                            ),
                            total_budget=float(sc.get("total_budget", 100000)),
                            industry=sc.get("industry", "Technology"),
                            roles=(
                                ", ".join(_sc_roles)
                                if isinstance(_sc_roles, list)
                                else str(_sc_roles)
                            ),
                            locations=(
                                ", ".join(_sc_locs)
                                if isinstance(_sc_locs, list)
                                else str(_sc_locs)
                            ),
                        )
                        simulated.append(sim_result)
                result = compare_scenarios(simulated)
                # LLM insights for Scenario Comparison
                if isinstance(result, dict):
                    result["ai_insights"] = _generate_product_insights(
                        "Budget Scenario Comparison",
                        {
                            k: result.get(k)
                            for k in ("winner", "comparison_matrix", "summary")
                            if result.get(k) is not None
                        },
                        context=f"Comparing {len(simulated)} budget scenarios",
                    )
                self._send_json(result)
            except Exception as e:
                logger.error("Simulator compare error: %s", e, exc_info=True)
                self._send_json({"error": _internal_error_msg()})
        elif path == "/api/simulator/export/excel":
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
                body = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body)
                from budget_simulator import export_scenario_excel

                excel_bytes = export_scenario_excel(
                    data.get("scenarios") or [],
                    data.get("comparison", {}),
                    data.get("client_name", "Simulation"),
                )
                self.send_response(200)
                self.send_header(
                    "Content-Type",
                    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                )
                self.send_header(
                    "Content-Disposition", "attachment; filename=budget_simulation.xlsx"
                )
                self.send_header("Content-Length", str(len(excel_bytes)))
                self.end_headers()
                self.wfile.write(excel_bytes)
            except Exception as e:
                logger.error("Simulator Excel export error: %s", e, exc_info=True)
                self.send_response(500)
                self.send_header("Content-Type", "application/json")
                self.end_headers()
                self.wfile.write(
                    _error_response("Excel export failed", "SERVER_ERROR", 500)[0]
                )
        elif path == "/api/simulator/export/ppt":
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
                body = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body)
                from budget_simulator import export_scenario_ppt

                ppt_bytes = export_scenario_ppt(
                    data.get("scenarios") or [],
                    data.get("comparison", {}),
                    data.get("client_name", "Simulation"),
                )
                self.send_response(200)
                self.send_header(
                    "Content-Type",
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
                self.send_header(
                    "Content-Disposition", "attachment; filename=budget_simulation.pptx"
                )
                self.send_header("Content-Length", str(len(ppt_bytes)))
                self.end_headers()
                self.wfile.write(ppt_bytes)
            except Exception as e:
                logger.error("Simulator PPT export error: %s", e, exc_info=True)
                self.send_response(500)
                self.send_header("Content-Type", "application/json")
                self.end_headers()
                self.wfile.write(
                    _error_response("PPT export failed", "SERVER_ERROR", 500)[0]
                )
        # NOTE: /api/nova/export and /api/deliver now handled by routes/export.py
        # NOTE: /api/competitive/* now handled by routes/competitive.py
        # ── API Portal: All POST routes ──
        elif path.startswith("/api/portal/"):
            if not self._check_admin_auth():
                self._send_json({"error": "Unauthorized"}, status_code=401)
                return
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
                body = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body)
                from api_portal import handle_portal_api

                headers = {k: v for k, v in self.headers.items()}
                status_code, result = handle_portal_api(
                    path.replace("/api/", ""), "POST", data, headers
                )
                self.send_response(status_code)
                self.send_header("Content-Type", "application/json")
                self.end_headers()
                self.wfile.write(json.dumps(result, default=str).encode())
            except Exception as e:
                logger.error("API Portal POST error: %s", e, exc_info=True)
                self._send_json({"error": _internal_error_msg(), "status": "error"})
        # ── Social & Search Media Plan: Generate ──
        elif path == "/api/social-plan/generate":
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
                body = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body)
                from social_plan import generate_social_media_plan

                result = generate_social_media_plan(
                    role=data.get("role") or "",
                    location=data.get("location") or "",
                    industry=data.get("industry", "general_entry_level"),
                    budget=data.get("budget", 50000),
                    goals=data.get("goals", ["job_applications"]),
                    duration_weeks=data.get("duration_weeks", 4),
                )

                # LLM insights for Social Plan
                if isinstance(result, dict):
                    result["ai_insights"] = _generate_product_insights(
                        "Social & Search Media Plan",
                        {
                            k: result.get(k)
                            for k in (
                                "platform_allocations",
                                "weekly_schedule",
                                "budget_breakdown",
                                "projected_reach",
                            )
                            if result.get(k) is not None
                        },
                        context=f"Role: {data.get('role') or ''}, "
                        f"Budget: ${data.get('budget', 50000):,.0f}, "
                        f"Duration: {data.get('duration_weeks', 4)} weeks",
                    )

                # S72: Firecrawl platform ad specs enrichment removed (module deleted).
                # Static specs in data/platform_ad_specs.json cover the same surface;
                # Magic UI / 21st.dev component generation also picks up specs at
                # render time, so the social plan still ships with correct dimensions.

                self._send_json(result)
            except Exception as e:
                logger.error("Social plan generation error: %s", e, exc_info=True)
                self._send_json({"success": False, "error": _internal_error_msg()})
        # ── Social & Search Media Plan: Excel Download ──
        elif path == "/api/social-plan/download/excel":
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
                body = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body)
                from social_plan import generate_social_plan_excel

                excel_bytes = generate_social_plan_excel(plan=data)
                self.send_response(200)
                self.send_header(
                    "Content-Type",
                    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                )
                self.send_header(
                    "Content-Disposition", "attachment; filename=social_media_plan.xlsx"
                )
                self.send_header("Content-Length", str(len(excel_bytes)))
                self.end_headers()
                self.wfile.write(excel_bytes)
            except Exception as e:
                logger.error("Social plan Excel error: %s", e, exc_info=True)
                self.send_response(500)
                self.send_header("Content-Type", "application/json")
                self.end_headers()
                self.wfile.write(
                    _error_response("Excel export failed", "SERVER_ERROR", 500)[0]
                )
        # ── Social & Search Media Plan: PPT Download ──
        elif path == "/api/social-plan/download/ppt":
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
                body = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body)
                from social_plan import generate_social_plan_ppt

                ppt_bytes = generate_social_plan_ppt(plan=data)
                self.send_response(200)
                self.send_header(
                    "Content-Type",
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
                self.send_header(
                    "Content-Disposition", "attachment; filename=social_media_plan.pptx"
                )
                self.send_header("Content-Length", str(len(ppt_bytes)))
                self.end_headers()
                self.wfile.write(ppt_bytes)
            except Exception as e:
                logger.error("Social plan PPT error: %s", e, exc_info=True)
                self.send_response(500)
                self.send_header("Content-Type", "application/json")
                self.end_headers()
                self.wfile.write(
                    _error_response("PPT export failed", "SERVER_ERROR", 500)[0]
                )
        # ── Talent Heatmap: Analyze ──
        elif path == "/api/talent-heatmap/analyze":
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
                body = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body)
                from talent_heatmap import run_heatmap_analysis

                result = run_heatmap_analysis(
                    role=data.get("role") or "",
                    industry=data.get("industry", "general_entry_level"),
                    locations=data.get("locations"),
                    budget=data.get("budget") or 0,
                    num_hires=data.get("num_hires", 5),
                )

                # Enrich with Supabase salary data
                if _supabase_data_available and isinstance(result, dict):
                    try:
                        _thm_role = data.get("role") or ""
                        _thm_locations = data.get("locations") or []
                        if isinstance(_thm_locations, str):
                            _thm_locations = [_thm_locations]
                        _thm_salary_data: list = []
                        for _thm_loc in _thm_locations[:10]:
                            _thm_loc_str = (
                                _thm_loc if isinstance(_thm_loc, str) else str(_thm_loc)
                            )
                            _sb_salaries = get_salary_data(
                                role=_thm_role, location=_thm_loc_str
                            )
                            if _sb_salaries:
                                _thm_salary_data.extend(_sb_salaries)
                        if _thm_salary_data:
                            result["supabase_salary_data"] = _thm_salary_data
                    except (urllib.error.URLError, OSError, ValueError) as sb_err:
                        logger.error(
                            f"Supabase salary enrichment for heatmap failed: {sb_err}",
                            exc_info=True,
                        )

                # ── API Integrations: Talent Heatmap enrichment ──
                if _api_integrations_available and isinstance(result, dict):
                    _thm_live_data: dict = {}
                    _thm_role_str = data.get("role") or ""
                    _thm_locs = data.get("locations") or []
                    if isinstance(_thm_locs, str):
                        _thm_locs = [_thm_locs]
                    # Adzuna job counts per location
                    if _api_adzuna and _thm_role_str:
                        try:
                            _thm_job_counts: dict = {}
                            for _loc in _thm_locs[:10]:
                                _loc_str = _loc if isinstance(_loc, str) else str(_loc)
                                _count = _api_adzuna.get_job_count(_thm_role_str, "us")
                                if _count:
                                    _thm_job_counts[_loc_str] = _count
                            if _thm_job_counts:
                                _thm_live_data["adzuna_job_counts"] = _thm_job_counts
                                logger.info(
                                    "Enriched /api/talent-heatmap with adzuna job_count data"
                                )
                        except (
                            urllib.error.URLError,
                            OSError,
                            ValueError,
                            TypeError,
                        ) as _ae:
                            logger.error(
                                "Adzuna enrichment for talent-heatmap failed: %s",
                                _ae,
                                exc_info=True,
                            )
                    # Census population by state
                    if _api_census:
                        try:
                            _census_pop = _api_census.get_population_by_state()
                            if _census_pop:
                                _thm_live_data["census_population"] = _census_pop
                                logger.info(
                                    "Enriched /api/talent-heatmap with census population data"
                                )
                        except (
                            urllib.error.URLError,
                            OSError,
                            ValueError,
                            TypeError,
                        ) as _ce:
                            logger.error(
                                "Census enrichment for talent-heatmap failed: %s",
                                _ce,
                                exc_info=True,
                            )
                    # BEA regional employment
                    if _api_bea:
                        try:
                            _bea_emp = _api_bea.get_regional_employment()
                            if _bea_emp:
                                _thm_live_data["bea_regional_employment"] = _bea_emp
                                logger.info(
                                    "Enriched /api/talent-heatmap with bea employment data"
                                )
                        except (
                            urllib.error.URLError,
                            OSError,
                            ValueError,
                            TypeError,
                        ) as _be:
                            logger.error(
                                "BEA enrichment for talent-heatmap failed: %s",
                                _be,
                                exc_info=True,
                            )
                    if _thm_live_data:
                        result["live_market_data"] = _thm_live_data

                # ── JobSpy: Real job density from 5 sources ──
                if (
                    _jobspy_available
                    and _jobspy_scrape_jobs
                    and isinstance(result, dict)
                ):
                    try:
                        _thm_js_role = data.get("role") or ""
                        _thm_js_locs = data.get("locations") or []
                        if isinstance(_thm_js_locs, str):
                            _thm_js_locs = [_thm_js_locs]
                        if _thm_js_role:
                            _thm_js_density: dict = {}
                            for _jsloc in _thm_js_locs[:5]:
                                _jsloc_str = (
                                    _jsloc if isinstance(_jsloc, str) else str(_jsloc)
                                )
                                _js_jobs = _jobspy_scrape_jobs(
                                    _thm_js_role, _jsloc_str, results_wanted=10
                                )
                                if _js_jobs is not None:
                                    _thm_js_density[_jsloc_str] = len(_js_jobs)
                            if _thm_js_density:
                                _live = result.get("live_market_data") or {}
                                _live["jobspy_posting_density"] = _thm_js_density
                                result["live_market_data"] = _live
                                logger.info(
                                    "Enriched /api/talent-heatmap with jobspy density data"
                                )
                    except (ValueError, TypeError, KeyError, OSError) as _jse:
                        logger.error(
                            "JobSpy enrichment for talent-heatmap failed: %s",
                            _jse,
                            exc_info=True,
                        )

                # LLM insights for Talent Heatmap
                if isinstance(result, dict):
                    result["ai_insights"] = _generate_product_insights(
                        "Talent Heatmap",
                        {
                            k: result.get(k)
                            for k in (
                                "heatmap_data",
                                "top_locations",
                                "cost_analysis",
                                "competition_index",
                            )
                            if result.get(k) is not None
                        },
                        context=f"Role: {data.get('role') or ''}, "
                        f"Industry: {data.get('industry', 'general_entry_level')}",
                    )

                # Enrich with Census Bureau workforce demographics
                try:
                    from api_enrichment import fetch_census_workforce_data

                    locations_list = data.get("locations") or []
                    if isinstance(locations_list, str):
                        locations_list = [locations_list]
                    census_demographics: dict = {}
                    for loc in locations_list[:10]:
                        loc_str = loc if isinstance(loc, str) else str(loc)
                        census = fetch_census_workforce_data(loc_str)
                        if census:
                            census_demographics[loc_str] = census
                    if census_demographics:
                        result["census_demographics"] = census_demographics
                except Exception as census_err:
                    logger.error(
                        "Census enrichment for heatmap failed: %s",
                        census_err,
                        exc_info=True,
                    )

                # S72: Firecrawl live job density removed (module deleted).
                # Adzuna + Jooble + USAJobs (in api_integrations) provide density
                # signals through the talent-heatmap baseline path; this Firecrawl
                # overlay was the supplementary live source.

                self._send_json(result)
            except Exception as e:
                logger.error("Talent heatmap analysis error: %s", e, exc_info=True)
                self._send_json({"error": _internal_error_msg(), "status": "error"})
        # ── Talent Heatmap: Excel Download ──
        elif path == "/api/talent-heatmap/download/excel":
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
                body = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body)
                from talent_heatmap import generate_heatmap_excel

                excel_bytes = generate_heatmap_excel(analysis=data)
                self.send_response(200)
                self.send_header(
                    "Content-Type",
                    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                )
                self.send_header(
                    "Content-Disposition", "attachment; filename=talent_heatmap.xlsx"
                )
                self.send_header("Content-Length", str(len(excel_bytes)))
                self.end_headers()
                self.wfile.write(excel_bytes)
            except Exception as e:
                logger.error("Talent heatmap Excel error: %s", e, exc_info=True)
                self.send_response(500)
                self.send_header("Content-Type", "application/json")
                self.end_headers()
                self.wfile.write(
                    _error_response("Excel export failed", "SERVER_ERROR", 500)[0]
                )
        # ── Talent Heatmap: PPT Download ──
        elif path == "/api/talent-heatmap/download/ppt":
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
                body = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body)
                from talent_heatmap import generate_heatmap_ppt

                ppt_bytes = generate_heatmap_ppt(analysis=data)
                self.send_response(200)
                self.send_header(
                    "Content-Type",
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
                self.send_header(
                    "Content-Disposition", "attachment; filename=talent_heatmap.pptx"
                )
                self.send_header("Content-Length", str(len(ppt_bytes)))
                self.end_headers()
                self.wfile.write(ppt_bytes)
            except Exception as e:
                logger.error("Talent heatmap PPT error: %s", e, exc_info=True)
                self.send_response(500)
                self.send_header("Content-Type", "application/json")
                self.end_headers()
                self.wfile.write(
                    _error_response("PPT export failed", "SERVER_ERROR", 500)[0]
                )
        # ── SkillTarget: Analyze / Excel / PPT ──
        elif path.startswith("/api/skill-target/"):
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
                body = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body)
                from skill_target import handle_skill_target_request

                result = handle_skill_target_request(path, "POST", data)
                # LLM insights for Skill Target
                if isinstance(result, dict) and "status_code" not in result:
                    result["ai_insights"] = _generate_product_insights(
                        "Skill Target Intelligence",
                        {
                            k: result.get(k)
                            for k in (
                                "skill_gaps",
                                "talent_pool",
                                "recommended_channels",
                                "difficulty_score",
                            )
                            if result.get(k) is not None
                        },
                        context=f"Skill targeting analysis",
                    )
                elif isinstance(result, dict) and result.get("status_code") == 200:
                    # Inject into JSON body of structured response
                    try:
                        _st_body = result.get("body") or ""
                        _st_parsed = (
                            json.loads(_st_body)
                            if isinstance(_st_body, str)
                            else _st_body
                        )
                        if isinstance(_st_parsed, dict):
                            _st_parsed["ai_insights"] = _generate_product_insights(
                                "Skill Target Intelligence",
                                {
                                    k: _st_parsed.get(k)
                                    for k in (
                                        "skill_gaps",
                                        "talent_pool",
                                        "recommended_channels",
                                        "difficulty_score",
                                    )
                                    if _st_parsed.get(k) is not None
                                },
                                context=f"Skill targeting analysis",
                            )
                            result["body"] = json.dumps(_st_parsed)
                    except (json.JSONDecodeError, TypeError, AttributeError):
                        pass  # Non-JSON body, skip insights
                # Unwrap structured response from skill_target.py
                if (
                    isinstance(result, dict)
                    and "status_code" in result
                    and "content_type" in result
                    and "body" in result
                ):
                    sc = result["status_code"]
                    ct = result["content_type"]
                    resp_body = result["body"]
                    if ct == "application/json":
                        # Parse JSON body string and send as JSON with proper status
                        parsed = (
                            json.loads(resp_body)
                            if isinstance(resp_body, str)
                            else resp_body
                        )
                        self.send_response(sc)
                        self.send_header("Content-Type", "application/json")
                        self.end_headers()
                        self.wfile.write(json.dumps(parsed).encode())
                    elif ct in (
                        "application/octet-stream",
                        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    ):
                        # Binary file download
                        raw = (
                            resp_body
                            if isinstance(resp_body, bytes)
                            else resp_body.encode()
                        )
                        if "/excel" in path:
                            fn = "skill_target_report.xlsx"
                        else:
                            fn = "skill_target_report.pptx"
                        self.send_response(sc)
                        self.send_header("Content-Type", ct)
                        self.send_header(
                            "Content-Disposition", f"attachment; filename={fn}"
                        )
                        self.send_header("Content-Length", str(len(raw)))
                        self.end_headers()
                        self.wfile.write(raw)
                    else:
                        # Other content types -- send as-is
                        raw = (
                            resp_body
                            if isinstance(resp_body, bytes)
                            else (
                                resp_body.encode()
                                if isinstance(resp_body, str)
                                else json.dumps(resp_body).encode()
                            )
                        )
                        self.send_response(sc)
                        self.send_header("Content-Type", ct)
                        self.send_header("Content-Length", str(len(raw)))
                        self.end_headers()
                        self.wfile.write(raw)
                elif isinstance(result, bytes):
                    # Legacy: raw binary file download (Excel or PPT)
                    if "/excel" in path:
                        ct = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                        fn = "skill_target_report.xlsx"
                    else:
                        ct = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        fn = "skill_target_report.pptx"
                    self.send_response(200)
                    self.send_header("Content-Type", ct)
                    self.send_header(
                        "Content-Disposition", f"attachment; filename={fn}"
                    )
                    self.send_header("Content-Length", str(len(result)))
                    self.end_headers()
                    self.wfile.write(result)
                else:
                    self._send_json(result)
            except Exception as e:
                logger.error("SkillTarget error: %s", e, exc_info=True)
                self._send_json(
                    {"error": "Analysis failed. Please try again.", "status": "error"}
                )
        # ── Market Intel Reports: Generate / Excel / PPT ──
        elif path.startswith("/api/market-intel/"):
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
                body = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body)
                from market_intel_reports import handle_market_intel_request

                result = handle_market_intel_request(path, "POST", data)
                if isinstance(result, bytes):
                    if "/excel" in path:
                        ct = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                        fn = "market_intel_report.xlsx"
                    else:
                        ct = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        fn = "market_intel_report.pptx"
                    self.send_response(200)
                    self.send_header("Content-Type", ct)
                    self.send_header(
                        "Content-Disposition", f"attachment; filename={fn}"
                    )
                    self.send_header("Content-Length", str(len(result)))
                    self.end_headers()
                    self.wfile.write(result)
                else:
                    self._send_json(result)
            except Exception as e:
                logger.error("Market Intel error: %s", e, exc_info=True)
                self._send_json(
                    {
                        "error": "Report generation failed. Please try again.",
                        "status": "error",
                    }
                )
        # ── Quick Plan Builder ──
        elif path == "/api/quick-plan":
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
                body = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body)
                from quick_plan import generate_quick_plan

                result = generate_quick_plan(
                    role=data.get("role") or "",
                    location=data.get("location") or "",
                    budget=data.get("budget", 50000),
                    industry=data.get("industry") or "",
                )
                self._send_json(result)
            except Exception as e:
                logger.error("Quick plan error: %s", e, exc_info=True)
                self._send_json({"success": False, "error": _internal_error_msg()})
        # ── Recruitment Advertising Audit: Full Audit ──
        elif path == "/api/audit/analyze":
            try:
                content_type = self.headers.get("Content-Type") or ""
                if "multipart/form-data" in content_type:
                    import cgi

                    form = cgi.FieldStorage(
                        fp=self.rfile,
                        headers=self.headers,
                        environ={
                            "REQUEST_METHOD": "POST",
                            "CONTENT_TYPE": content_type,
                        },
                    )
                    file_item = form["file"]
                    file_bytes = file_item.file.read()
                    filename = file_item.filename
                    industry = form.getvalue("industry", "general_entry_level")
                    _roles = form.getvalue("roles", "")
                    _locations = form.getvalue("locations", "")
                    client_name = form.getvalue("client_name", "Client")
                    roles = (
                        [r.strip() for r in _roles.split(",") if r.strip()]
                        if _roles
                        else None
                    )
                    locations = (
                        [l.strip() for l in _locations.split(",") if l.strip()]
                        if _locations
                        else None
                    )
                else:
                    content_len = int(self.headers.get("Content-Length") or 0)
                    body = self.rfile.read(content_len) if content_len > 0 else b"{}"
                    data = json.loads(body)
                    import base64

                    file_bytes = base64.b64decode(data.get("file_base64") or "")
                    filename = data.get("filename", "plan.xlsx")
                    industry = data.get("industry", "general_entry_level")
                    roles = data.get("roles")
                    locations = data.get("locations")
                    client_name = data.get("client_name", "Client")
                from audit_tool import run_full_audit

                result = run_full_audit(
                    file_bytes=file_bytes,
                    filename=filename,
                    industry=industry,
                    roles=roles,
                    locations=locations,
                    client_name=client_name,
                )
                # LLM insights for Audit Tool
                if isinstance(result, dict):
                    result["ai_insights"] = _generate_product_insights(
                        "Recruitment Advertising Audit",
                        {
                            k: result.get(k)
                            for k in (
                                "overall_score",
                                "findings",
                                "spend_efficiency",
                                "channel_audit",
                            )
                            if result.get(k) is not None
                        },
                        context=f"Client: {client_name}, Industry: {industry}",
                    )
                self._send_json(result)
            except Exception as e:
                logger.error("Audit analysis error: %s", e, exc_info=True)
                self._send_json({"error": _internal_error_msg(), "success": False})
        # ── Audit: Excel Download ──
        elif path == "/api/audit/download/excel":
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
                body = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body)
                from audit_tool import generate_audit_excel

                excel_bytes = generate_audit_excel(
                    report_data=data.get("report_data", data),
                    client_name=data.get("client_name", "Client"),
                )
                self.send_response(200)
                self.send_header(
                    "Content-Type",
                    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                )
                self.send_header(
                    "Content-Disposition", "attachment; filename=recruitment_audit.xlsx"
                )
                self.send_header("Content-Length", str(len(excel_bytes)))
                self.end_headers()
                self.wfile.write(excel_bytes)
            except Exception as e:
                logger.error("Audit Excel error: %s", e, exc_info=True)
                self.send_response(500)
                self.send_header("Content-Type", "application/json")
                self.end_headers()
                self.wfile.write(
                    _error_response("Excel export failed", "SERVER_ERROR", 500)[0]
                )
        # ── Audit: PPT Download ──
        elif path == "/api/audit/download/ppt":
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
                body = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body)
                from audit_tool import generate_audit_ppt

                ppt_bytes = generate_audit_ppt(
                    report_data=data.get("report_data", data),
                    client_name=data.get("client_name", "Client"),
                )
                self.send_response(200)
                self.send_header(
                    "Content-Type",
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
                self.send_header(
                    "Content-Disposition", "attachment; filename=recruitment_audit.pptx"
                )
                self.send_header("Content-Length", str(len(ppt_bytes)))
                self.end_headers()
                self.wfile.write(ppt_bytes)
            except Exception as e:
                logger.error("Audit PPT error: %s", e, exc_info=True)
                self.send_response(500)
                self.send_header("Content-Type", "application/json")
                self.end_headers()
                self.wfile.write(
                    _error_response("PPT export failed", "SERVER_ERROR", 500)[0]
                )
        # ── HireSignal: Analyze ──
        elif path == "/api/hire-signal/analyze":
            try:
                content_type = self.headers.get("Content-Type") or ""
                if "multipart/form-data" in content_type:
                    import cgi

                    form = cgi.FieldStorage(
                        fp=self.rfile,
                        headers=self.headers,
                        environ={
                            "REQUEST_METHOD": "POST",
                            "CONTENT_TYPE": content_type,
                        },
                    )
                    file_item = form["file"]
                    file_bytes = file_item.file.read()
                    filename = file_item.filename
                    industry = form.getvalue("industry", "general_entry_level")
                    client_name = form.getvalue("client_name", "Client")
                    from hire_signal import run_full_signal_analysis

                    result = run_full_signal_analysis(
                        file_bytes=file_bytes,
                        filename=filename,
                        industry=industry,
                        client_name=client_name,
                    )
                else:
                    content_len = int(self.headers.get("Content-Length") or 0)
                    body = self.rfile.read(content_len) if content_len > 0 else b"{}"
                    data = json.loads(body)
                    from hire_signal import run_full_signal_analysis

                    # Validate manual_data entries are dicts, not strings
                    _raw_manual = data.get("manual_data")
                    if isinstance(_raw_manual, list):
                        _raw_manual = [
                            item if isinstance(item, dict) else {"source": str(item)}
                            for item in _raw_manual
                        ]

                    result = run_full_signal_analysis(
                        manual_data=_raw_manual,
                        industry=data.get("industry") or "general_entry_level",
                        client_name=data.get("client_name") or "Client",
                    )

                # LLM insights for HireSignal
                if isinstance(result, dict):
                    result["ai_insights"] = _generate_product_insights(
                        "HireSignal Intelligence",
                        {
                            k: result.get(k)
                            for k in (
                                "signal_score",
                                "risk_factors",
                                "recommendations",
                                "market_signals",
                                "channel_analysis",
                            )
                            if result.get(k) is not None
                        },
                        context=f"Industry: {(data.get('industry') or '') if isinstance(data, dict) else ''}",
                    )

                # Enrich with Google Trends hiring data
                if isinstance(result, dict):
                    try:
                        from api_enrichment import fetch_hiring_trends

                        hs_industry = (
                            industry
                            if "multipart/form-data"
                            in (self.headers.get("Content-Type") or "")
                            else (data.get("industry") or "")
                        )
                        if hs_industry:
                            trends = fetch_hiring_trends(hs_industry)
                            if trends:
                                result["google_trends"] = trends
                    except Exception as trends_err:
                        logger.error(
                            "Google Trends enrichment for HireSignal failed: %s",
                            trends_err,
                            exc_info=True,
                        )

                    # S72: Firecrawl job posting volume removed (module deleted).
                    # Adzuna and BLS JOLTS (in api_integrations) provide posting
                    # volume in the HireSignal baseline path.

                    # Enrich with Supabase supply repository data
                    if _supabase_data_available:
                        try:
                            _hs_industry = (
                                industry
                                if "multipart/form-data"
                                in (self.headers.get("Content-Type") or "")
                                else (data.get("industry") or "")
                            )
                            _hs_supply = get_supply_repository(category=_hs_industry)
                            if _hs_supply:
                                result["supabase_supply_data"] = _hs_supply
                        except (urllib.error.URLError, OSError, ValueError) as sb_err:
                            logger.error(
                                f"Supabase supply repository fetch for HireSignal failed: {sb_err}",
                                exc_info=True,
                            )

                    # ── API Integrations: HireSignal enrichment ──
                    if _api_integrations_available:
                        _hs_api_data: dict = {}
                        _hs_role_str = (
                            result.get("role") or result.get("job_title") or ""
                        )
                        _hs_loc_str = result.get("location") or ""
                        # Adzuna current job postings
                        if _api_adzuna and _hs_role_str:
                            try:
                                _hs_jobs = _api_adzuna.search_jobs(
                                    _hs_role_str, "us", results_per_page=5
                                )
                                if _hs_jobs:
                                    _hs_api_data["adzuna_current_postings"] = _hs_jobs
                                    logger.info(
                                        "Enriched /api/hire-signal with adzuna job postings"
                                    )
                            except (
                                urllib.error.URLError,
                                OSError,
                                ValueError,
                                TypeError,
                            ) as _ae:
                                logger.error(
                                    "Adzuna enrichment for hire-signal failed: %s",
                                    _ae,
                                    exc_info=True,
                                )
                        # O*NET skill requirements
                        if _api_onet and _hs_role_str:
                            try:
                                _onet_occs = _api_onet.search_occupations(_hs_role_str)
                                if _onet_occs and len(_onet_occs) > 0:
                                    _top_soc = _onet_occs[0].get("code") or ""
                                    if _top_soc:
                                        _onet_skills = _api_onet.get_skills(_top_soc)
                                        if _onet_skills:
                                            _hs_api_data["onet_skills"] = _onet_skills[
                                                :10
                                            ]
                                            logger.info(
                                                "Enriched /api/hire-signal with onet skills data"
                                            )
                            except (
                                urllib.error.URLError,
                                OSError,
                                ValueError,
                                TypeError,
                            ) as _oe:
                                logger.error(
                                    "O*NET enrichment for hire-signal failed: %s",
                                    _oe,
                                    exc_info=True,
                                )
                        # USAJobs government sector comparison
                        if _api_usajobs and _hs_role_str:
                            try:
                                _usa_jobs = _api_usajobs.search_jobs(
                                    _hs_role_str, _hs_loc_str, results_per_page=5
                                )
                                if _usa_jobs:
                                    _hs_api_data["usajobs_federal_postings"] = _usa_jobs
                                    logger.info(
                                        "Enriched /api/hire-signal with usajobs data"
                                    )
                            except (
                                urllib.error.URLError,
                                OSError,
                                ValueError,
                                TypeError,
                            ) as _ue:
                                logger.error(
                                    "USAJobs enrichment for hire-signal failed: %s",
                                    _ue,
                                    exc_info=True,
                                )
                        if _hs_api_data:
                            result["live_market_data"] = _hs_api_data

                    # ── JobSpy: Multi-source job data for HireSignal ──
                    if _jobspy_available and _jobspy_market_stats:
                        try:
                            _hs_js_role = (
                                result.get("role") or result.get("job_title") or ""
                            )
                            if _hs_js_role:
                                _hs_js_stats = _jobspy_market_stats(str(_hs_js_role))
                                if _hs_js_stats:
                                    _hs_live = result.get("live_market_data") or {}
                                    _hs_live["jobspy_market_stats"] = _hs_js_stats
                                    result["live_market_data"] = _hs_live
                                    logger.info(
                                        "Enriched /api/hire-signal with jobspy market stats"
                                    )
                        except (ValueError, TypeError, KeyError, OSError) as _jse:
                            logger.error(
                                "JobSpy enrichment for hire-signal failed: %s",
                                _jse,
                                exc_info=True,
                            )

                self._send_json(result)
            except Exception as e:
                logger.error("HireSignal analyze error: %s", e, exc_info=True)
                self._send_json({"error": _internal_error_msg(), "success": False})
        # ── HireSignal: Sample Data ──
        elif path == "/api/hire-signal/sample-data":
            try:
                from hire_signal import generate_sample_data, run_full_signal_analysis

                sample = generate_sample_data()
                content_len = int(self.headers.get("Content-Length") or 0)
                body = self.rfile.read(content_len) if content_len > 0 else b"{}"
                try:
                    data = json.loads(body)
                except Exception:
                    data = {}
                result = run_full_signal_analysis(
                    manual_data=sample,
                    industry=data.get("industry", "tech_engineering"),
                    client_name=data.get("client_name", "Demo Company"),
                )
                self._send_json(result)
            except Exception as e:
                logger.error("HireSignal sample data error: %s", e, exc_info=True)
                self._send_json({"error": _internal_error_msg(), "success": False})
        # ── HireSignal: Excel Download ──
        elif path == "/api/hire-signal/download/excel":
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
                body = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body)
                from hire_signal import generate_signal_excel

                excel_bytes = generate_signal_excel(
                    report=data.get("report_data", data),
                    client_name=data.get("client_name", "Client"),
                )
                self.send_response(200)
                self.send_header(
                    "Content-Type",
                    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                )
                self.send_header(
                    "Content-Disposition",
                    "attachment; filename=hire_signal_report.xlsx",
                )
                self.send_header("Content-Length", str(len(excel_bytes)))
                self.end_headers()
                self.wfile.write(excel_bytes)
            except Exception as e:
                logger.error("HireSignal Excel error: %s", e, exc_info=True)
                self.send_response(500)
                self.send_header("Content-Type", "application/json")
                self.end_headers()
                self.wfile.write(
                    _error_response("Excel export failed", "SERVER_ERROR", 500)[0]
                )
        # ── HireSignal: PPT Download ──
        elif path == "/api/hire-signal/download/ppt":
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
                body = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body)
                from hire_signal import generate_signal_ppt

                ppt_bytes = generate_signal_ppt(
                    report=data.get("report_data", data),
                    client_name=data.get("client_name", "Client"),
                )
                self.send_response(200)
                self.send_header(
                    "Content-Type",
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
                self.send_header(
                    "Content-Disposition",
                    "attachment; filename=hire_signal_report.pptx",
                )
                self.send_header("Content-Length", str(len(ppt_bytes)))
                self.end_headers()
                self.wfile.write(ppt_bytes)
            except Exception as e:
                logger.error("HireSignal PPT error: %s", e, exc_info=True)
                self.send_response(500)
                self.send_header("Content-Type", "application/json")
                self.end_headers()
                self.wfile.write(
                    _error_response("PPT export failed", "SERVER_ERROR", 500)[0]
                )
        # ── Market Pulse POST APIs ──
        elif path.startswith("/api/pulse/"):
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
                body_raw = self.rfile.read(content_len) if content_len > 0 else b"{}"
                body_data = json.loads(body_raw)
                from market_pulse import handle_pulse_api

                status_code, result = handle_pulse_api(
                    path, method="POST", body=body_data
                )

                # Enrich with Supabase market trends
                if (
                    _supabase_data_available
                    and status_code == 200
                    and isinstance(result, dict)
                ):
                    try:
                        _mp_category = (
                            body_data.get("category") or body_data.get("sector") or ""
                        )
                        _mp_trends = get_market_trends(category=_mp_category)
                        if _mp_trends:
                            result["supabase_market_trends"] = _mp_trends
                    except (urllib.error.URLError, OSError, ValueError) as sb_err:
                        logger.error(
                            f"Supabase market trends fetch failed: {sb_err}",
                            exc_info=True,
                        )

                # Enrich with Google Trends hiring data
                if status_code == 200 and isinstance(result, dict):
                    try:
                        from api_enrichment import fetch_hiring_trends

                        industry = (
                            body_data.get("industry") or body_data.get("sector") or ""
                        )
                        if industry:
                            trends = fetch_hiring_trends(industry)
                            if trends:
                                result["google_trends"] = trends
                    except Exception as trends_err:
                        logger.error(
                            "Google Trends enrichment for pulse failed: %s",
                            trends_err,
                            exc_info=True,
                        )

                if status_code == 200:
                    self._send_json(result)
                else:
                    self.send_response(status_code)
                    self.send_header("Content-Type", "application/json")
                    self.end_headers()
                    self.wfile.write(json.dumps(result).encode())
            except Exception as e:
                logger.error("Market Pulse POST error: %s", e, exc_info=True)
                self._send_json({"ok": False, "error": _internal_error_msg()})
        # ── ApplyFlow API ──
        elif path == "/api/applyflow":
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
                body_raw = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body_raw)
                from applyflow import handle_applyflow_request

                result = handle_applyflow_request(data)
                self._send_json(result)
            except Exception as e:
                logger.error("ApplyFlow error: %s", e, exc_info=True)
                self._send_json(
                    {
                        "error": "Request failed",
                        "response": "Something went wrong. Please try again.",
                    }
                )
        # ── Post-Campaign Analysis: Excel Download ──
        elif path == "/api/post-campaign/download/excel":
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
                body_raw = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body_raw)

                wb = Workbook()

                # -- Sheet 1: Campaign Summary --
                ws_summary = wb.active
                ws_summary.title = "Campaign Summary"
                header_font = Font(name="Calibri", bold=True, size=12, color="FFFFFF")
                header_fill = PatternFill(
                    start_color="6366F1", end_color="6366F1", fill_type="solid"
                )
                subheader_font = Font(name="Calibri", bold=True, size=10)
                normal_font = Font(name="Calibri", size=10)
                border_thin = Border(
                    left=Side(style="thin", color="D9D9D9"),
                    right=Side(style="thin", color="D9D9D9"),
                    top=Side(style="thin", color="D9D9D9"),
                    bottom=Side(style="thin", color="D9D9D9"),
                )

                # Title row
                ws_summary.merge_cells("A1:B1")
                ws_summary["A1"] = "Post-Campaign Analysis"
                ws_summary["A1"].font = Font(
                    name="Calibri", bold=True, size=16, color="6366F1"
                )

                summary_fields = [
                    ("Campaign Name", data.get("campaign_name") or ""),
                    ("Client Name", data.get("client_name") or ""),
                    ("Start Date", data.get("start_date") or ""),
                    ("End Date", data.get("end_date") or ""),
                    ("Total Budget", data.get("total_budget") or 0),
                    ("Industry", data.get("industry") or ""),
                    ("Benchmark CPA", data.get("benchmark_cpa") or 0),
                ]
                for i, (label, val) in enumerate(summary_fields, start=3):
                    ws_summary.cell(row=i, column=1, value=label).font = subheader_font
                    ws_summary.cell(row=i, column=2, value=val).font = normal_font

                ws_summary.column_dimensions["A"].width = 22
                ws_summary.column_dimensions["B"].width = 35

                # -- Sheet 2: Channel Performance --
                ws_perf = wb.create_sheet("Channel Performance")
                perf_headers = [
                    "Channel",
                    "Spend ($)",
                    "Impressions",
                    "Clicks",
                    "Applications",
                    "Hires",
                    "CTR (%)",
                    "Apply Rate (%)",
                    "CPA ($)",
                    "CPH ($)",
                    "Grade",
                ]
                for col_idx, hdr in enumerate(perf_headers, start=1):
                    cell = ws_perf.cell(row=1, column=col_idx, value=hdr)
                    cell.font = header_font
                    cell.fill = header_fill
                    cell.alignment = Alignment(horizontal="center")
                    cell.border = border_thin

                channels = data.get("channels") or []
                if isinstance(channels, list) and len(channels) > 100:
                    self._send_error(
                        "Channels list exceeds 100 items limit", "VALIDATION_ERROR", 400
                    )
                    return
                total_spend = 0
                total_apps = 0
                total_hires = 0
                for row_idx, ch in enumerate(channels, start=2):
                    spend = ch.get("spend") or 0
                    impressions = ch.get("impressions") or 0
                    clicks = ch.get("clicks") or 0
                    applications = ch.get("applications") or 0
                    hires = ch.get("hires") or 0
                    ctr = ch.get("ctr") or 0
                    apply_rate = ch.get("applyRate") or 0
                    cpa = ch.get("cpa") or 0
                    cph = ch.get("cph") or 0
                    grade = ch.get("grade", "--")
                    total_spend += spend
                    total_apps += applications
                    total_hires += hires

                    vals = [
                        ch.get("channel") or "",
                        spend,
                        impressions,
                        clicks,
                        applications,
                        hires,
                        round(ctr, 2),
                        round(apply_rate, 2),
                        round(cpa, 2),
                        round(cph, 2),
                        grade,
                    ]
                    for col_idx, val in enumerate(vals, start=1):
                        cell = ws_perf.cell(row=row_idx, column=col_idx, value=val)
                        cell.font = normal_font
                        cell.border = border_thin
                        cell.alignment = Alignment(horizontal="center")

                # Totals row
                tot_row = len(channels) + 2
                ws_perf.cell(row=tot_row, column=1, value="TOTAL").font = Font(
                    name="Calibri", bold=True, size=10
                )
                ws_perf.cell(row=tot_row, column=2, value=total_spend).font = Font(
                    name="Calibri", bold=True, size=10
                )
                ws_perf.cell(row=tot_row, column=5, value=total_apps).font = Font(
                    name="Calibri", bold=True, size=10
                )
                ws_perf.cell(row=tot_row, column=6, value=total_hires).font = Font(
                    name="Calibri", bold=True, size=10
                )
                overall_cpa = total_spend / total_apps if total_apps > 0 else 0
                overall_cph = total_spend / total_hires if total_hires > 0 else 0
                ws_perf.cell(
                    row=tot_row, column=9, value=round(overall_cpa, 2)
                ).font = Font(name="Calibri", bold=True, size=10)
                ws_perf.cell(
                    row=tot_row, column=10, value=round(overall_cph, 2)
                ).font = Font(name="Calibri", bold=True, size=10)
                for col_idx in range(1, 12):
                    ws_perf.cell(row=tot_row, column=col_idx).border = border_thin

                # Auto-width columns
                for col_idx in range(1, 12):
                    ws_perf.column_dimensions[get_column_letter(col_idx)].width = 16

                # Save to buffer
                buf = io.BytesIO()
                wb.save(buf)
                excel_bytes = buf.getvalue()

                self.send_response(200)
                self.send_header(
                    "Content-Type",
                    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                )
                self.send_header(
                    "Content-Disposition",
                    "attachment; filename=post_campaign_analysis.xlsx",
                )
                self.send_header("Content-Length", str(len(excel_bytes)))
                self.end_headers()
                self.wfile.write(excel_bytes)
            except Exception as e:
                logger.error("Post-campaign Excel error: %s", e, exc_info=True)
                self.send_response(500)
                self.send_header("Content-Type", "application/json")
                self.end_headers()
                self.wfile.write(
                    _error_response(
                        f"Excel generation failed: {e}", "SERVER_ERROR", 500
                    )[0]
                )
        # ── ComplianceGuard: Analyze ──
        elif path == "/api/compliance/analyze":
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
                body_raw = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body_raw)
                result = _analyze_compliance(data)

                # Return 400 if required fields are missing
                if (
                    isinstance(result, dict)
                    and result.get("error")
                    and not data.get("text")
                ):
                    self._send_json(result, status_code=400)
                    return

                # Enrich with Supabase compliance rules
                if _supabase_data_available and isinstance(result, dict):
                    try:
                        _cg_location = data.get("location") or ""
                        _cg_rules = get_compliance_rules(jurisdiction=_cg_location)
                        if _cg_rules:
                            result["supabase_compliance_rules"] = _cg_rules
                    except (urllib.error.URLError, OSError, ValueError) as sb_err:
                        logger.error(
                            f"Supabase compliance rules fetch failed: {sb_err}",
                            exc_info=True,
                        )

                # S72: Firecrawl compliance law updates removed (module deleted).
                # Static compliance rules in Supabase compliance_rules table cover
                # the same surface (already loaded above via get_compliance_rules).

                self._send_json(result)
                # ── PostHog: Track compliance audit ──
                self._ph_track(
                    "compliance_audit_run",
                    {
                        "endpoint": "/api/compliance/analyze",
                    },
                )
            except json.JSONDecodeError:
                self._send_error("Invalid JSON", "VALIDATION_ERROR", 400)
            except Exception as e:
                logger.error("Compliance analyze error: %s", e, exc_info=True)
                self.send_response(500)
                self.send_header("Content-Type", "application/json")
                self.end_headers()
                self.wfile.write(
                    _error_response(
                        f"Compliance analysis failed: {e}", "SERVER_ERROR", 500
                    )[0]
                )
        # ── ComplianceGuard: Audit (enhanced endpoint) ──
        elif path == "/api/complianceguard/audit":
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
                body_raw = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body_raw)
                result = _audit_complianceguard(data)

                if (
                    isinstance(result, dict)
                    and result.get("error")
                    and not data.get("job_description")
                ):
                    self._send_json(result, status_code=400)
                    return

                self._send_json(result)
                # ── PostHog: Track compliance audit ──
                self._ph_track(
                    "compliance_audit_run",
                    {
                        "endpoint": "/api/complianceguard/audit",
                    },
                )
            except json.JSONDecodeError:
                self._send_error("Invalid JSON", "VALIDATION_ERROR", 400)
            except Exception as e:
                logger.error(f"ComplianceGuard audit error: {e}", exc_info=True)
                self.send_response(500)
                self.send_header("Content-Type", "application/json")
                self.end_headers()
                self.wfile.write(
                    _error_response(
                        f"ComplianceGuard audit failed: {e}", "SERVER_ERROR", 500
                    )[0]
                )
        # ── CreativeAI: Generate ──
        elif path == "/api/creative/generate":
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
                body_raw = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body_raw)
                result = _generate_creative_ads(data)
                # Return 400 if required fields are missing
                if isinstance(result, dict) and result.get("error"):
                    self._send_json(result, status_code=400)
                else:
                    self._send_json(result)
            except json.JSONDecodeError:
                self._send_error("Invalid JSON", "VALIDATION_ERROR", 400)
            except Exception as e:
                logger.error("Creative generate error: %s", e, exc_info=True)
                self.send_response(500)
                self.send_header("Content-Type", "application/json")
                self.end_headers()
                self.wfile.write(
                    _error_response(
                        f"Creative generation failed: {e}", "SERVER_ERROR", 500
                    )[0]
                )
        # ── Post-Campaign Analysis: AI Summary ──
        elif path == "/api/post-campaign/summary":
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
                body_raw = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body_raw)
                result = _generate_post_campaign_summary(data)
                self._send_json(result)
            except json.JSONDecodeError:
                self._send_error("Invalid JSON", "VALIDATION_ERROR", 400)
            except Exception as e:
                logger.error("Post-campaign summary error: %s", e, exc_info=True)
                self.send_response(500)
                self.send_header("Content-Type", "application/json")
                self.end_headers()
                self.wfile.write(
                    _error_response(
                        f"Post-campaign summary failed: {e}", "SERVER_ERROR", 500
                    )[0]
                )
        # ── ROI Calculator: Calculate ──
        elif path == "/api/roi/calculate":
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
                body_raw = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body_raw)
                result = _calculate_roi(data)
                # Return 400 if required fields are invalid
                if isinstance(result, dict) and result.get("error"):
                    self._send_json(result, status_code=400)
                else:
                    self._send_json(result)
            except json.JSONDecodeError:
                self._send_error("Invalid JSON", "VALIDATION_ERROR", 400)
            except Exception as e:
                logger.error("ROI calculate error: %s", e, exc_info=True)
                self.send_response(500)
                self.send_header("Content-Type", "application/json")
                self.end_headers()
                self.wfile.write(
                    _error_response(
                        f"ROI calculation failed: {e}", "SERVER_ERROR", 500
                    )[0]
                )
        # ── ROI Calculator: Excel Download ──
        elif path == "/api/roi/download/excel":
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
                body_raw = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body_raw)
                excel_bytes = _generate_roi_excel(data)
                self.send_response(200)
                self.send_header(
                    "Content-Type",
                    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                )
                self.send_header(
                    "Content-Disposition",
                    "attachment; filename=roi_projection.xlsx",
                )
                self.send_header("Content-Length", str(len(excel_bytes)))
                self.end_headers()
                self.wfile.write(excel_bytes)
            except Exception as e:
                logger.error("ROI Excel error: %s", e, exc_info=True)
                self.send_response(500)
                self.send_header("Content-Type", "application/json")
                self.end_headers()
                self.wfile.write(
                    _error_response(
                        f"ROI Excel generation failed: {e}", "SERVER_ERROR", 500
                    )[0]
                )
        # ── A/B Test Lab: Generate ──
        elif path == "/api/ab-test/generate":
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
                body_raw = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body_raw)
                result = _generate_ab_test(data)
                self._send_json(result)
            except json.JSONDecodeError:
                self._send_error("Invalid JSON", "VALIDATION_ERROR", 400)
            except Exception as e:
                logger.error("A/B test generate error: %s", e, exc_info=True)
                self.send_response(500)
                self.send_header("Content-Type", "application/json")
                self.end_headers()
                self.wfile.write(
                    _error_response(
                        f"A/B test generation failed: {e}", "SERVER_ERROR", 500
                    )[0]
                )
        # ── Post-Campaign Analysis: PPT Download ──
        elif path == "/api/post-campaign/download/ppt":
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
                body_raw = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body_raw)
                ppt_bytes = _generate_post_campaign_ppt(data)
                self.send_response(200)
                self.send_header(
                    "Content-Type",
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
                self.send_header(
                    "Content-Disposition",
                    "attachment; filename=post_campaign_analysis.pptx",
                )
                self.send_header("Content-Length", str(len(ppt_bytes)))
                self.end_headers()
                self.wfile.write(ppt_bytes)
            except Exception as e:
                logger.error("Post-campaign PPT error: %s", e, exc_info=True)
                self.send_response(500)
                self.send_header("Content-Type", "application/json")
                self.end_headers()
                self.wfile.write(
                    _error_response(f"PPT generation failed: {e}", "SERVER_ERROR", 500)[
                        0
                    ]
                )
        # NOTE: /api/vendor-iq/live-pricing and /api/payscale-sync/salary now handled by routes/pricing.py
        # NOTE: /api/report/html, /api/export/sheets, /api/export/status now handled by routes/export.py
        # ── Sentry Webhook (POST) ──
        # M3: Signature verification is performed inside _handle_sentry_webhook
        # via validate_sentry_signature(). Requires SENTRY_WEBHOOK_SECRET env var.
        # If not set, the handler rejects the webhook (returns 401).
        elif path == "/api/sentry/webhook":
            try:
                if not _sentry_integration_available:
                    self._send_json(
                        {"error": "sentry_integration module not available"},
                        status_code=503,
                    )
                    return

                try:
                    content_len = int(self.headers.get("Content-Length") or 0)
                except (ValueError, TypeError):
                    content_len = 0
                if content_len > 1 * 1024 * 1024:  # 1MB limit
                    self._send_json({"error": "Request too large"}, status_code=413)
                    return

                body = self.rfile.read(content_len) if content_len > 0 else b""
                signature = self.headers.get("Sentry-Hook-Signature") or ""
                resource = self.headers.get("Sentry-Hook-Resource") or ""

                status_code, result = _handle_sentry_webhook(
                    body=body,
                    signature=signature,
                    resource_header=resource,
                )
                self._send_json(result, status_code=status_code)

            except Exception as e:
                logger.error("Sentry webhook error: %s", e, exc_info=True)
                self._send_json({"error": _internal_error_msg()}, status_code=500)

        # ── ElevenLabs: Text-to-Speech ──
        elif path == "/api/elevenlabs/tts":
            if not _elevenlabs_available:
                self._send_json(
                    {"error": "ElevenLabs integration not available"}, status_code=503
                )
                return
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
                body_raw = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body_raw)
                text = str(data.get("text") or "")
                if not text.strip():
                    self._send_json(
                        {"error": "Missing required field: text"}, status_code=400
                    )
                    return
                voice_id = data.get("voice_id") or None
                model_id = data.get("model_id") or "eleven_flash_v2_5"
                audio = _elevenlabs_tts(text, voice_id=voice_id, model_id=model_id)
                if audio:
                    self.send_response(200)
                    self.send_header("Content-Type", "audio/mpeg")
                    self.send_header("Content-Length", str(len(audio)))
                    cors_origin = self._get_cors_origin()
                    if cors_origin:
                        self.send_header("Access-Control-Allow-Origin", cors_origin)
                        self.send_header("Access-Control-Allow-Credentials", "true")
                    self.end_headers()
                    self.wfile.write(audio)
                else:
                    self._send_json({"error": "TTS generation failed"}, status_code=502)
            except json.JSONDecodeError:
                self._send_json({"error": "Invalid JSON"}, status_code=400)
            except Exception as e:
                logger.error("ElevenLabs TTS error: %s", e, exc_info=True)
                self._send_json({"error": "Internal TTS error"}, status_code=500)

        # ── ElevenLabs: Speech-to-Text ──
        elif path == "/api/elevenlabs/stt":
            if not _elevenlabs_available:
                self._send_json(
                    {"error": "ElevenLabs integration not available"}, status_code=503
                )
                return
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
                _STT_MAX_BYTES = 25 * 1024 * 1024  # 25 MB
                if content_len == 0:
                    self._send_json(
                        {"error": "No audio data provided"}, status_code=400
                    )
                    return
                if content_len > _STT_MAX_BYTES:
                    self._send_json(
                        {"error": "Audio file too large (max 25MB)"}, status_code=413
                    )
                    return
                audio_data = self.rfile.read(content_len)
                content_type = self.headers.get("Content-Type") or ""
                # If JSON body with base64, decode it
                if "application/json" in content_type:
                    import base64

                    data = json.loads(audio_data)
                    audio_b64 = data.get("audio") or ""
                    language = data.get("language") or "en"
                    if not audio_b64:
                        self._send_json(
                            {"error": "Missing required field: audio (base64)"},
                            status_code=400,
                        )
                        return
                    audio_data = base64.b64decode(audio_b64)
                else:
                    language = "en"
                transcript = _elevenlabs_stt(audio_data, language=language)
                if transcript is not None:
                    self._send_json({"text": transcript, "language": language})
                else:
                    self._send_json(
                        {"error": "Speech-to-text transcription failed"},
                        status_code=502,
                    )
            except json.JSONDecodeError:
                self._send_json({"error": "Invalid JSON"}, status_code=400)
            except Exception as e:
                logger.error("ElevenLabs STT error: %s", e, exc_info=True)
                self._send_json({"error": "Internal STT error"}, status_code=500)

        # ── ElevenLabs: Sound Effects ──
        elif path == "/api/elevenlabs/sfx":
            if not _elevenlabs_available:
                self._send_json(
                    {"error": "ElevenLabs integration not available"}, status_code=503
                )
                return
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
                body_raw = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body_raw)
                description = data.get("description") or ""
                if not description.strip():
                    self._send_json(
                        {"error": "Missing required field: description"},
                        status_code=400,
                    )
                    return
                duration = float(data.get("duration_seconds") or 5.0)
                audio = _elevenlabs_sfx(description, duration_seconds=duration)
                if audio:
                    self.send_response(200)
                    self.send_header("Content-Type", "audio/mpeg")
                    self.send_header("Content-Length", str(len(audio)))
                    cors_origin = self._get_cors_origin()
                    if cors_origin:
                        self.send_header("Access-Control-Allow-Origin", cors_origin)
                        self.send_header("Access-Control-Allow-Credentials", "true")
                    self.end_headers()
                    self.wfile.write(audio)
                else:
                    self._send_json(
                        {"error": "Sound effect generation failed"}, status_code=502
                    )
            except json.JSONDecodeError:
                self._send_json({"error": "Invalid JSON"}, status_code=400)
            except (ValueError, TypeError) as e:
                self._send_json({"error": f"Invalid parameter: {e}"}, status_code=400)
            except Exception as e:
                logger.error("ElevenLabs SFX error: %s", e, exc_info=True)
                self._send_json(
                    {"error": "Internal sound effect error"}, status_code=500
                )

        # ── ElevenLabs: Ad Voiceover (CreativeAI) ──
        elif path == "/api/creative/voiceover":
            if not _elevenlabs_available:
                self._send_json(
                    {"error": "ElevenLabs integration not available"}, status_code=503
                )
                return
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
                body_raw = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body_raw)
                ad_copy = data.get("ad_copy") or data.get("text") or ""
                if not ad_copy.strip():
                    self._send_json(
                        {"error": "Missing required field: ad_copy"}, status_code=400
                    )
                    return
                tone = data.get("tone") or "professional"
                audio = _elevenlabs_voiceover(ad_copy, tone=tone)
                if audio:
                    self.send_response(200)
                    self.send_header("Content-Type", "audio/mpeg")
                    self.send_header("Content-Length", str(len(audio)))
                    cors_origin = self._get_cors_origin()
                    if cors_origin:
                        self.send_header("Access-Control-Allow-Origin", cors_origin)
                        self.send_header("Access-Control-Allow-Credentials", "true")
                    self.end_headers()
                    self.wfile.write(audio)
                else:
                    self._send_json(
                        {"error": "Ad voiceover generation failed"}, status_code=502
                    )
            except json.JSONDecodeError:
                self._send_json({"error": "Invalid JSON"}, status_code=400)
            except Exception as e:
                logger.error("ElevenLabs voiceover error: %s", e, exc_info=True)
                self._send_json({"error": "Internal voiceover error"}, status_code=500)

        # ── Proactive Intelligence Insights (POST) ──
        elif path == "/api/insights/read":
            try:
                from nova_proactive import mark_insight_read

                content_len = int(self.headers.get("Content-Length") or 0)
                body = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body)
                insight_id = data.get("id") or ""
                result = mark_insight_read(insight_id)
                self._send_json({"ok": result})
            except Exception as e:
                logger.error("Insight mark-read error: %s", e, exc_info=True)
                self._send_json({"error": _internal_error_msg()}, status_code=500)

        elif path == "/api/insights/dismiss":
            try:
                from nova_proactive import dismiss_insight

                content_len = int(self.headers.get("Content-Length") or 0)
                body = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body)
                insight_id = data.get("id") or ""
                result = dismiss_insight(insight_id)
                self._send_json({"ok": result})
            except Exception as e:
                logger.error("Insight dismiss error: %s", e, exc_info=True)
                self._send_json({"error": _internal_error_msg()}, status_code=500)

        # ── Attribution Report (POST /api/attribution/report) ──
        elif path == "/api/attribution/report":
            try:
                if not _attribution_available:
                    self._send_json(
                        {"error": "Attribution module not available"},
                        status_code=503,
                    )
                    return
                content_len = int(self.headers.get("Content-Length") or 0)
                body = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body)
                if (
                    not data.get("total_budget")
                    and not data.get("budget")
                    and not data.get("channels")
                ):
                    self._send_json(
                        {
                            "error": "Request must include total_budget or channels with spend data"
                        },
                        status_code=400,
                    )
                    return
                report = _attribution_report(data)
                self._send_json(report)
            except json.JSONDecodeError:
                self._send_json({"error": "Invalid JSON body"}, status_code=400)
            except Exception as e:
                logger.error("Attribution report error: %s", e, exc_info=True)
                self._send_json(
                    {"error": f"Attribution report failed: {e}"}, status_code=500
                )

        # ── Plan Outcome Prediction Model ──
        elif path == "/api/predict":
            try:
                from prediction_model import predict_outcomes

                content_len = int(self.headers.get("Content-Length") or 0)
                body = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body)
                if not data:
                    self._send_json(
                        {"error": "Request body must contain plan data"},
                        status_code=400,
                    )
                    return
                result = predict_outcomes(data)
                self._send_json(result)
            except json.JSONDecodeError:
                self._send_json({"error": "Invalid JSON body"}, status_code=400)
            except ImportError:
                self._send_json(
                    {"error": "Prediction model not available"}, status_code=503
                )
            except Exception as e:
                logger.error("Prediction error: %s", e, exc_info=True)
                self._send_json({"error": f"Prediction failed: {e}"}, status_code=500)

        elif path == "/api/predict/grade":
            try:
                from prediction_model import grade_plan

                content_len = int(self.headers.get("Content-Length") or 0)
                body = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body)
                if not data:
                    self._send_json(
                        {"error": "Request body must contain plan data"},
                        status_code=400,
                    )
                    return
                result = grade_plan(data)
                self._send_json(result)
            except json.JSONDecodeError:
                self._send_json({"error": "Invalid JSON body"}, status_code=400)
            except ImportError:
                self._send_json(
                    {"error": "Prediction model not available"}, status_code=503
                )
            except Exception as e:
                logger.error("Plan grading error: %s", e, exc_info=True)
                self._send_json({"error": f"Plan grading failed: {e}"}, status_code=500)

        elif path == "/api/predict/compare":
            try:
                from prediction_model import compare_plans

                content_len = int(self.headers.get("Content-Length") or 0)
                body = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body)
                plan_a = data.get("plan_a")
                plan_b = data.get("plan_b")
                if not plan_a or not plan_b:
                    self._send_json(
                        {"error": "Request must include both 'plan_a' and 'plan_b'"},
                        status_code=400,
                    )
                    return
                result = compare_plans(plan_a, plan_b)
                self._send_json(result)
            except json.JSONDecodeError:
                self._send_json({"error": "Invalid JSON body"}, status_code=400)
            except ImportError:
                self._send_json(
                    {"error": "Prediction model not available"}, status_code=503
                )
            except Exception as e:
                logger.error("Plan comparison error: %s", e, exc_info=True)
                self._send_json(
                    {"error": f"Plan comparison failed: {e}"}, status_code=500
                )

        # ── Outcome Pipeline: record outcome event ──
        elif path == "/api/outcomes/record":
            if not self._check_rate_limit():
                self._send_error("Rate limit exceeded", "RATE_LIMITED", 429)
                return
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
                if content_len <= 0 or content_len > 1024 * 1024:
                    _sz_status = 400 if content_len <= 0 else 413
                    _sz_msg = (
                        "Empty request body"
                        if content_len <= 0
                        else "Request too large"
                    )
                    self._send_error(_sz_msg, "VALIDATION_ERROR", _sz_status)
                    return
                body = self.rfile.read(content_len)
                data = json.loads(body)
                plan_id = (data.get("plan_id") or "").strip()
                stage = (data.get("stage") or "").strip()
                if not plan_id or not stage:
                    self._send_error(
                        "Missing required fields: plan_id, stage",
                        "VALIDATION_ERROR",
                        400,
                    )
                    return

                from outcome_pipeline import record_outcome, update_prediction_model

                result = record_outcome(
                    plan_id=plan_id,
                    stage=stage,
                    data=data.get("data"),
                    role_family=data.get("role_family") or "",
                    predicted=data.get("predicted"),
                )

                # Auto-trigger model update on terminal stages
                if stage in ("hires_made", "offers_extended"):
                    try:
                        update_prediction_model()
                    except Exception as model_err:
                        logger.error(
                            "Model update after outcome recording failed: %s",
                            model_err,
                            exc_info=True,
                        )

                self._send_json(result, status_code=201)
            except json.JSONDecodeError:
                self._send_error("Invalid JSON body", "VALIDATION_ERROR", 400)
            except Exception as e:
                logger.error("Outcome record error: %s", e, exc_info=True)
                self._send_json(
                    {"error": f"Failed to record outcome: {e}"}, status_code=500
                )

        # ── Calendar Sync MCP: POST /api/calendar/events ──
        elif path == "/api/calendar/events":
            if not _calendar_sync_available:
                self._send_json(
                    {"error": "Calendar sync not available"}, status_code=503
                )
                return
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
                body_raw = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body_raw)
                title = data.get("title") or ""
                date = data.get("date") or ""
                details = data.get("details") or ""
                duration = int(data.get("duration_hours") or 1)
                attendees = data.get("attendees")
                if not title or not date:
                    self._send_json(
                        {"error": "title and date are required"}, status_code=400
                    )
                    return
                event = _cal_create_event(
                    title=title,
                    date=date,
                    details=details,
                    duration_hours=duration,
                    attendees=attendees,
                )
                if event:
                    self._send_json({"success": True, "event": event})
                else:
                    self._send_json(
                        {"success": False, "error": "Failed to create event"},
                        status_code=500,
                    )
            except json.JSONDecodeError:
                self._send_json({"error": "Invalid JSON"}, status_code=400)
            except Exception as exc:
                logger.error("Calendar create event error: %s", exc, exc_info=True)
                if _slack_alerts_available:
                    _slack_error_alert(exc, context="POST /api/calendar/events")
                self._send_json({"error": "Calendar service error"}, status_code=500)

        # ── Chroma RAG MCP: POST /api/nova/rag/search ──
        elif path == "/api/nova/rag/search":
            if not _chroma_rag_available:
                self._send_json(
                    {"error": "ChromaDB RAG not available", "results": []},
                    status_code=503,
                )
                return
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
                body_raw = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body_raw)
                query = data.get("query") or data.get("q") or ""
                top_k = int(data.get("top_k") or 5)
                top_k = max(1, min(top_k, 20))
                if not query:
                    self._send_json(
                        {"error": "query is required", "results": []},
                        status_code=400,
                    )
                    return
                results = _chroma_search(query, top_k=top_k)
                self._send_json(
                    {"results": results, "count": len(results), "query": query}
                )
            except json.JSONDecodeError:
                self._send_json(
                    {"error": "Invalid JSON", "results": []}, status_code=400
                )
            except Exception as exc:
                logger.error("Chroma RAG search POST error: %s", exc, exc_info=True)
                if _slack_alerts_available:
                    _slack_error_alert(exc, context="POST /api/nova/rag/search")
                self._send_json(
                    {"error": "RAG search service error", "results": []},
                    status_code=500,
                )

        # ── Slack Alerts MCP: POST /api/slack/alerts/send ──
        elif path == "/api/slack/alerts/send":
            if not self._check_admin_auth():
                self._send_json({"error": "Unauthorized"}, status_code=401)
                return
            if not _slack_alerts_available:
                self._send_json(
                    {"error": "Slack alerts not available"}, status_code=503
                )
                return
            try:
                content_len = int(self.headers.get("Content-Length") or 0)
                body_raw = self.rfile.read(content_len) if content_len > 0 else b"{}"
                data = json.loads(body_raw)
                message = data.get("message") or ""
                severity = data.get("severity") or "info"
                channel = data.get("channel")
                if not message:
                    self._send_json({"error": "message is required"}, status_code=400)
                    return
                success = _slack_send_alert(
                    channel=channel, message=message, severity=severity
                )
                self._send_json({"success": success})
            except json.JSONDecodeError:
                self._send_json({"error": "Invalid JSON"}, status_code=400)
            except Exception as exc:
                logger.error("Slack alert send error: %s", exc, exc_info=True)
                self._send_json({"error": _internal_error_msg()}, status_code=500)

        # ── AI Co-Pilot: Inline suggestions for media plan generator ──
        # NOTE: /api/copilot/suggest now handled by routes/copilot.py
        # NOTE: /api/plan/share and /api/plan/feedback now handled by routes/campaign.py
        else:
            self.send_error(404)

    def do_HEAD(self) -> None:
        """Handle HEAD requests for monitoring tools and uptime checks.

        Delegates to do_GET; BaseHTTPRequestHandler automatically suppresses
        the response body for HEAD requests when the command attribute is 'HEAD'.
        """
        self.do_GET()

    def do_OPTIONS(self) -> None:
        """Handle CORS preflight requests."""
        self.send_response(204)
        cors_origin = self._get_cors_origin()
        if cors_origin:
            self.send_header("Access-Control-Allow-Origin", cors_origin)
            self.send_header("Access-Control-Allow-Credentials", "true")
        self.send_header("Access-Control-Allow-Methods", "GET, HEAD, POST, OPTIONS")
        self.send_header(
            "Access-Control-Allow-Headers",
            "Content-Type, Authorization, X-Async, X-CSRF-Token, X-Idempotency-Key",
        )
        self.send_header("Access-Control-Max-Age", "86400")
        self.end_headers()

    def _send_compressed_response(
        self, data: bytes, content_type: str, status: int = 200
    ) -> None:
        """Send a gzip-compressed HTTP response if the client accepts it.

        Falls back to uncompressed response for small payloads (<1KB)
        or clients that don't advertise gzip support.
        """
        accept_encoding = self.headers.get("Accept-Encoding") or ""
        if "gzip" in accept_encoding and len(data) > 1024:
            compressed = gzip_module.compress(data, compresslevel=6)
            self.send_response(status)
            self.send_header("Content-Type", content_type)
            self.send_header("Content-Encoding", "gzip")
            self.send_header("Content-Length", str(len(compressed)))
            self.send_header("Vary", "Accept-Encoding")
            self.end_headers()
            self.wfile.write(compressed)
        else:
            self.send_response(status)
            self.send_header("Content-Type", content_type)
            self.send_header("Content-Length", str(len(data)))
            self.end_headers()
            self.wfile.write(data)

    # Nova chat widget snippet injected into all HTML template pages.
    # Excluded from nova.html (dedicated full-page chat), legal/terms pages,
    # and the API portal.
    _NOVA_WIDGET_SNIPPET = (
        b"\n<!-- Nova AI Chat Widget (auto-injected) -->"
        b'\n<script src="/static/nova-chat.js?v=3.5.3"></script>'
        b'\n<script>(function(){if(typeof NovaChat!=="undefined")NovaChat.init({containerId:null})})();</script>\n'
    )
    _NOVA_WIDGET_SKIP_PAGES = {
        "nova.html",
        "terms.html",
        "privacy.html",
        "api-portal.html",
        "platform.html",
    }

    def _serve_fragment(self, fragment_name: str) -> None:
        """Serve an HTML fragment for the Nova Platform SPA shell.

        Reads the corresponding template file, strips the outer HTML/head/body
        wrapper, and returns only the inner content. Results are cached in memory
        for performance. Thread-safe via _fragment_cache_lock.
        """
        if not fragment_name or fragment_name not in _FRAGMENT_MAP:
            self.send_error(
                404, f"Fragment not found: {html_mod.escape(fragment_name)}"
            )
            return

        # Check cache first (thread-safe)
        with _fragment_cache_lock:
            cached = _fragment_cache.get(fragment_name)
        if cached is not None:
            body = cached.encode("utf-8")
            self.send_response(200)
            self.send_header("Content-Type", "text/html; charset=utf-8")
            self.send_header("Cache-Control", "private, max-age=300")
            cors_origin = self._get_cors_origin()
            if cors_origin:
                self.send_header("Access-Control-Allow-Origin", cors_origin)
            self.send_header("Content-Length", str(len(body)))
            self.end_headers()
            self.wfile.write(body)
            return

        # Read template file (use composer for split templates)
        template_file = _FRAGMENT_MAP[fragment_name]
        template_path = os.path.join(TEMPLATES_DIR, template_file)
        page_name = template_file.rsplit(".", 1)[0]  # "index.html" -> "index"

        raw_html = None
        try:
            from template_composer import get_composed_template

            composed = get_composed_template(page_name)
            if composed is not None:
                raw_html = composed.decode("utf-8")
        except ImportError:
            pass

        if raw_html is None:
            if not os.path.exists(template_path):
                self.send_error(
                    404, f"Template not found: {html_mod.escape(template_file)}"
                )
                return
            try:
                with open(template_path, "r", encoding="utf-8") as f:
                    raw_html = f.read()
            except OSError as e:
                logger.error(
                    "Fragment read error for %s: %s", fragment_name, e, exc_info=True
                )
                self.send_error(500, "Failed to read template")
                return

        # Strip outer HTML/head/body wrapper, keep inner content
        fragment_html = self._extract_fragment_content(raw_html)

        # Cache the result (thread-safe)
        with _fragment_cache_lock:
            _fragment_cache[fragment_name] = fragment_html

        body = fragment_html.encode("utf-8")
        self.send_response(200)
        self.send_header("Content-Type", "text/html; charset=utf-8")
        self.send_header("Cache-Control", "private, max-age=300")
        cors_origin = self._get_cors_origin()
        if cors_origin:
            self.send_header("Access-Control-Allow-Origin", cors_origin)
        self.send_header("Content-Length", str(len(body)))
        self.end_headers()
        self.wfile.write(body)

    @staticmethod
    def _extract_fragment_content(html: str) -> str:
        """Extract the inner content from a full HTML page for fragment loading.

        Strips <!DOCTYPE>, <html>, <head>, and <body> tags, returning only the
        content that was inside <body>. Preserves inline <style> and <script>
        blocks that appear within the body. The <head> content (meta tags, title,
        external links) is discarded since the platform shell provides those.
        """
        # Remove DOCTYPE
        content = re.sub(r"<!DOCTYPE[^>]*>", "", html, flags=re.IGNORECASE)

        # Extract <style> blocks from <head> so CSS is preserved in fragments
        head_styles = ""
        head_match = re.search(
            r"<head[^>]*>(.*?)</head>",
            content,
            flags=re.DOTALL | re.IGNORECASE,
        )
        if head_match:
            head_content = head_match.group(1)
            style_blocks = re.findall(
                r"<style[^>]*>.*?</style>",
                head_content,
                flags=re.DOTALL | re.IGNORECASE,
            )
            if style_blocks:
                head_styles = "\n".join(style_blocks)

        # Extract body content if <body> tags exist
        body_match = re.search(
            r"<body[^>]*>(.*)</body>",
            content,
            flags=re.DOTALL | re.IGNORECASE,
        )
        if body_match:
            content = body_match.group(1)
        else:
            # No <body> tag -- strip <html> and <head> if present
            content = re.sub(r"</?html[^>]*>", "", content, flags=re.IGNORECASE)
            head_match_2 = re.search(
                r"<head[^>]*>.*?</head>",
                content,
                flags=re.DOTALL | re.IGNORECASE,
            )
            if head_match_2:
                content = (
                    content[: head_match_2.start()] + content[head_match_2.end() :]
                )

        # Prepend head styles to body content so CSS works in fragments
        if head_styles:
            content = head_styles + "\n" + content

        # Remove any stray nova-chat.js widget injection (platform has its own)
        content = re.sub(
            r"<!--\s*Nova AI Chat Widget.*?</script>\s*",
            "",
            content,
            flags=re.DOTALL | re.IGNORECASE,
        )

        return content.strip()

    def _serve_file(self, filepath, content_type):
        try:
            # For HTML templates, check if a composed version exists (split partials)
            if content_type == "text/html" and filepath.startswith(TEMPLATES_DIR):
                basename = os.path.basename(filepath)
                page_name = basename.rsplit(".", 1)[0]  # "index.html" -> "index"
                try:
                    from template_composer import get_composed_template

                    composed = get_composed_template(page_name)
                    if composed is not None:
                        content = composed
                    else:
                        with open(filepath, "rb") as f:
                            content = f.read()
                except ImportError:
                    with open(filepath, "rb") as f:
                        content = f.read()
            else:
                with open(filepath, "rb") as f:
                    content = f.read()

            # Auto-inject Nova chat widget into HTML template pages
            # (skip pages that already include it or shouldn't have it)
            if (
                content_type == "text/html"
                and filepath.startswith(TEMPLATES_DIR)
                and os.path.basename(filepath) not in self._NOVA_WIDGET_SKIP_PAGES
                and b"nova-chat.js" not in content
            ):
                # Insert just before </body>
                close_body = content.rfind(b"</body>")
                if close_body != -1:
                    content = (
                        content[:close_body]
                        + self._NOVA_WIDGET_SNIPPET
                        + content[close_body:]
                    )

            accept_encoding = self.headers.get("Accept-Encoding") or ""
            if "gzip" in accept_encoding and len(content) > 1024:
                compressed = gzip_module.compress(content, compresslevel=6)
                self.send_response(200)
                self.send_header("Content-Type", f"{content_type}; charset=utf-8")
                self.send_header("Content-Encoding", "gzip")
                self.send_header("Content-Length", str(len(compressed)))
                self.send_header("Vary", "Accept-Encoding")
                self.send_header("Cache-Control", "no-cache, no-store, must-revalidate")
                self.end_headers()
                self.wfile.write(compressed)
            else:
                self.send_response(200)
                self.send_header("Content-Type", f"{content_type}; charset=utf-8")
                self.send_header("Content-Length", str(len(content)))
                self.send_header("Cache-Control", "no-cache, no-store, must-revalidate")
                self.end_headers()
                self.wfile.write(content)
        except FileNotFoundError:
            self.send_error(404)

    # ── Cached 404 page (loaded once on first use) ──
    _404_page_cache: Optional[bytes] = None

    def send_error(
        self, code: int, message: Optional[str] = None, explain: Optional[str] = None
    ) -> None:
        """Override to serve branded 404 page for HTML requests.

        For API requests (Accept: application/json or /api/ paths), falls back
        to the default JSON-style error. For browser requests hitting a 404,
        serves the branded templates/404.html page.
        """
        if code == 404:
            accept = self.headers.get("Accept") or ""
            path = getattr(self, "path", "") or ""
            # Serve HTML 404 for browser requests, not API calls
            if "application/json" not in accept and not path.startswith("/api/"):
                try:
                    if self.__class__._404_page_cache is None:
                        p = os.path.join(TEMPLATES_DIR, "404.html")
                        with open(p, "rb") as f:
                            self.__class__._404_page_cache = f.read()
                    body = self.__class__._404_page_cache
                    self.send_response(404)
                    self.send_header("Content-Type", "text/html; charset=utf-8")
                    self.send_header("Content-Length", str(len(body)))
                    self.send_header("Cache-Control", "no-cache")
                    self.end_headers()
                    self.wfile.write(body)
                    return
                except OSError:
                    pass  # Fall through to default error
        super().send_error(code, message, explain)

    def send_response(self, code: int, message=None) -> None:
        """Override to track the response status code for metrics."""
        self._response_status = code
        super().send_response(code, message)

    def _send_json(self, data: Any, status_code: int = 200) -> None:
        """Send a JSON response.

        Success responses (2xx) are sent as-is (no envelope wrapping).
        Error responses (4xx/5xx) passed through _send_json (legacy callers
        that pass {"error": ...} instead of using _send_error) are wrapped as:
            {success: false, data: null, error: <message>, code: "ERROR"}

        Already-enveloped responses (containing a top-level "success" key)
        are passed through unchanged.
        """
        if isinstance(data, dict) and "success" in data:
            body = json.dumps(data).encode("utf-8")
        elif status_code >= 400:
            # Legacy error path: callers passing {"error": "..."} directly
            err_msg = (
                data.get("error", "Unknown error")
                if isinstance(data, dict)
                else str(data)
            )
            envelope = {
                "success": False,
                "error": err_msg,
                "code": "ERROR",
                "data": None,
            }
            body = json.dumps(envelope).encode("utf-8")
        else:
            # S21: Send raw data -- all frontend code expects unwrapped responses
            body = json.dumps(data).encode("utf-8")
        self._send_compressed_response(body, "application/json", status=status_code)
        # Store in idempotency cache if applicable
        try:
            parsed_path = urlparse(self.path).path
            self._store_idempotency(parsed_path, body, status_code)
        except Exception:
            pass  # Non-critical -- never fail the response for caching


class ThreadedHTTPServer(ThreadingMixIn, HTTPServer):
    """Handle each request in a separate thread so /api/generate (20-60s) does
    not block health checks, static files, or other concurrent requests.
    daemon_threads ensures threads don't block process shutdown."""

    daemon_threads = True
    # Allow port reuse to avoid "Address already in use" on quick restarts
    allow_reuse_address = True
    # Cap pending connection backlog to prevent resource exhaustion
    request_queue_size = 50


if __name__ == "__main__":
    import signal

    # Production uses gunicorn via wsgi.py (see Procfile).
    # This __main__ block is for local development only.

    port = int(os.environ.get("PORT", sys.argv[1] if len(sys.argv) > 1 else 5001))
    server = ThreadedHTTPServer(("0.0.0.0", port), MediaPlanHandler)

    # ── Graceful shutdown on SIGTERM (Render.com sends SIGTERM) ──
    def _handle_signal(signum, frame):
        sig_name = signal.Signals(signum).name
        logger.info("Received %s -- initiating graceful shutdown", sig_name)
        if _shutdown:
            _shutdown.request_shutdown()
            _shutdown.wait_for_completion()
        # Flush PostHog events before shutdown
        if _posthog_available:
            try:
                _ph_shutdown()
            except Exception:
                pass
        server.shutdown()
        logger.info("Server shut down cleanly")

    signal.signal(signal.SIGTERM, _handle_signal)
    signal.signal(signal.SIGINT, _handle_signal)

    # ── Deferred startup: run heavy init AFTER server starts ──
    # This ensures the health check passes quickly while modules load in background
    def _deferred_startup():
        """Run all heavy initialization in a background thread."""
        import time as _t

        _t.sleep(2)  # Let server start accepting connections first
        logger.info("[deferred_startup] Beginning background initialization...")

        # Pre-compose split templates (fast -- do early)
        try:
            from template_composer import precompose_all

            composed = precompose_all()
            if composed:
                logger.info(
                    "[deferred_startup] Pre-composed %d templates: %s",
                    len(composed),
                    {k: f"{v:,}B" for k, v in composed.items()},
                )
        except ImportError:
            pass
        except Exception as tc_err:
            logger.error("[deferred_startup] Template pre-compose failed: %s", tc_err)

        # Pre-warm knowledge base
        try:
            kb = load_knowledge_base()
            logger.info(
                "[deferred_startup] Knowledge base pre-warmed: %d keys", len(kb)
            )
        except Exception as kb_err:
            logger.error(
                "[deferred_startup] Knowledge base pre-warm failed: %s", kb_err
            )

        # Build vector search index
        if _vector_search_available and _vector_index_kb:
            try:
                count = _vector_index_kb()
                logger.info(
                    "[deferred_startup] Vector search index built: %d documents", count
                )
            except Exception as ve:
                logger.error(
                    "[deferred_startup] Vector index build failed: %s",
                    ve,
                    exc_info=True,
                )

        # Data Refresh Pipeline
        try:
            from data_refresh import start_data_refresh

            start_data_refresh()
            logger.info("[deferred_startup] Data refresh pipeline started")
        except ImportError:
            logger.debug("[deferred_startup] data_refresh module not available")

        # Proactive Health Checker
        try:
            from sentry_integration import (
                start_proactive_health as _start_proactive_health,
            )

            _start_proactive_health()
            logger.info("[deferred_startup] Proactive health checker started")
        except ImportError:
            pass

        # Proactive Intelligence Engine
        try:
            from nova_proactive import start_proactive_engine

            start_proactive_engine()
            logger.info("[deferred_startup] Proactive intelligence engine started")
        except ImportError:
            pass

        # Feature Store Init
        try:
            from feature_store import get_feature_store

            get_feature_store().initialize()
            logger.info("[deferred_startup] Feature store initialized")
        except ImportError:
            pass
        except Exception as _fs_err:
            logger.error(
                "[deferred_startup] Feature store init failed: %s",
                _fs_err,
                exc_info=True,
            )

        # API Key Authentication Init
        global _auth_module_loaded
        try:
            from auth import init as _init_auth

            _init_auth()
            _auth_module_loaded = True
        except ImportError:
            _auth_module_loaded = False
            logger.critical(
                "SECURITY: auth module failed to import -- admin endpoints will be BLOCKED. "
                "Fix auth.py or set NOVA_AUTH_DISABLED=1 to explicitly run without auth."
            )

        # Pre-compute salary + demand data (runs after 10-min warm-up, then every 24h)
        try:
            from precompute import start_precompute_thread

            start_precompute_thread()
            logger.info(
                "[deferred_startup] Salary/demand pre-computation thread started"
            )
        except ImportError:
            logger.debug("[deferred_startup] precompute module not available")
        except Exception as _pc_err:
            logger.error(
                "[deferred_startup] Pre-compute thread start failed: %s",
                _pc_err,
                exc_info=True,
            )

        # Mark deploy warmup as complete (blue-green readiness gate)
        global _DEPLOY_WARMUP_COMPLETE
        _DEPLOY_WARMUP_COMPLETE = True
        logger.info(
            "[deferred_startup] Background initialization complete -- deploy warmup READY"
        )

    threading.Thread(
        target=_deferred_startup, daemon=True, name="deferred-startup"
    ).start()

    # ── Startup banner ──
    logger.info("=" * 60)
    logger.info("AI Media Planner v4.0.0")
    logger.info("Port: %d | PID: %d | Threads: daemon", port, os.getpid())
    logger.info("Health: http://localhost:%d/health", port)
    logger.info("Readiness: http://localhost:%d/ready", port)
    logger.info("Metrics: http://localhost:%d/api/metrics", port)
    logger.info(
        "PPTX: %s | API enrichment: %s",
        "available" if generate_pptx else "unavailable",
        "available" if enrich_data else "unavailable",
    )
    logger.info("Data Matrix: %s", "monitoring" if _data_matrix else "unavailable")
    logger.info("AutoQC: %s", "running" if _auto_qc else "unavailable")
    logger.info(
        "Data Enrichment: %s",
        "running" if _data_enrichment_available else "unavailable",
    )
    logger.info(
        "JobSpy: %s | Tavily: %s | VectorSearch: %s",
        "available" if _jobspy_available else "unavailable",
        "available" if _tavily_available else "unavailable",
        "available" if _vector_search_available else "unavailable",
    )
    logger.info("API Docs: http://localhost:%d/docs", port)
    logger.info("OpenAPI: http://localhost:%d/api/docs/openapi.json", port)
    logger.info("API Version: v1 | Async generation: enabled")
    logger.info("=" * 60)

    try:
        server.serve_forever()
    except KeyboardInterrupt:
        logger.info("KeyboardInterrupt received -- shutting down")
        server.shutdown()
    finally:
        logger.info("Server process exiting")
